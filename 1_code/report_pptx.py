# ===== BARUM report_pptx.py v31b — rep 기반 편집가능 보장진단서 PPT =====
# 설명지 PDF(build_report_pdf)와 동일 rep를 읽어 내용 일치 보장. CI/비CI 자동 분기.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# BARUM 브랜드 색 (설명지와 동일 톤: 딥 네이비 + 골드)
NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
GOLD  = RGBColor(0x9A, 0x7A, 0x12)
GREEN = RGBColor(0x1E, 0x7A, 0x45)
AMBER = RGBColor(0xB9, 0x54, 0x0B)
RED   = RGBColor(0xC0, 0x39, 0x2B)
GREY  = RGBColor(0x6A, 0x70, 0x7A)
LGREY = RGBColor(0xEE, 0xF0, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE  = RGBColor(0x1F, 0x5F, 0xA8)
DARK  = RGBColor(0x22, 0x28, 0x33)

SW, SH = Inches(13.333), Inches(7.5)   # 16:9

def _pct_color(p):
    return GREEN if p>=70 else (GOLD if p>=40 else RED)

def _box(slide, l, t, w, h, fill=None, line=None, line_w=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(0.75)
    return sp

def _txt(slide, l, t, w, h, text, size=12, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='맑은 고딕'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    for i, line in enumerate(str(text).split('\n')):
        para = p if i == 0 else tf.add_paragraph()
        para.alignment = align
        r = para.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb

def _header(slide, page_no, total, client, subtitle):
    # 상단 네이비 바
    bar = _box(slide, 0, 0, SW, Inches(1.05), fill=NAVY)
    _txt(slide, Inches(0.5), Inches(0.12), Inches(9), Inches(0.35),
         'MAKEONE · 보장분석 리포트', size=11, color=GOLD, bold=True)
    _txt(slide, Inches(0.5), Inches(0.42), Inches(9), Inches(0.55),
         f'{client} 고객님 보장 진단서', size=22, color=WHITE, bold=True)
    _txt(slide, Inches(11.2), Inches(0.18), Inches(1.9), Inches(0.4),
         str(page_no), size=26, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    _txt(slide, Inches(10.7), Inches(0.62), Inches(2.4), Inches(0.3),
         subtitle, size=10, color=RGBColor(0xC8,0xD0,0xDC), align=PP_ALIGN.RIGHT)
    # 하단 푸터
    _box(slide, 0, Inches(7.15), SW, Inches(0.35), fill=NAVY)
    _txt(slide, Inches(0.5), Inches(7.18), Inches(6), Inches(0.28),
         'MAKEONE 보장분석 자동화', size=9, color=RGBColor(0xC8,0xD0,0xDC))
    _txt(slide, Inches(7), Inches(7.18), Inches(5.8), Inches(0.28),
         f'{client} 고객님 · {page_no} / {total}', size=9,
         color=RGBColor(0xC8,0xD0,0xDC), align=PP_ALIGN.RIGHT)

def _slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    return s

def build_report_pptx(rep, out):
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    client = rep.get('client', '고객')
    ci = rep.get('ci', {}) or {}
    ci_present = bool(ci.get('present'))
    TOTAL = 8

    # ---------- P1 표지 ----------
    s = _slide(prs)
    _box(s, 0, 0, SW, SH, fill=NAVY)
    _txt(s, Inches(0.8), Inches(0.7), Inches(6), Inches(0.5),
         'MAKEONE', size=20, color=GOLD, bold=True)
    _txt(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.0),
         f'{client} 고객님', size=44, color=WHITE, bold=True)
    _txt(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.6),
         '보장 진단서', size=28, color=RGBColor(0xC8,0xD0,0xDC), bold=True)
    _box(s, Inches(0.85), Inches(4.5), Inches(4.2), Inches(0.05), fill=GOLD)
    _txt(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.5),
         f"보유 {rep.get('n_contract',0)}건 · 월 {rep.get('premium',0):,}원 · "
         f"갱신 {rep.get('renew',0)} / 비갱신 {rep.get('nonrenew',0)}",
         size=14, color=WHITE)
    _txt(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
         f"{rep.get('branch','')}  ·  {rep.get('manager','')} {rep.get('title','')}",
         size=11, color=RGBColor(0x9A,0xA4,0xB4))

    # ---------- P2 보장현황 ----------
    s = _slide(prs); _header(s, 2, TOTAL, client, '보장 현황')
    y = Inches(1.3)
    cov = rep.get('coverage', [])
    col_w = Inches(6.3); x_l = Inches(0.4); x_r = Inches(6.9)
    row_h = Inches(0.92)
    for i, c in enumerate(cov):
        col_x = x_l if i % 2 == 0 else x_r
        row_y = y + (i // 2) * row_h
        stc = {'full':GREEN,'part':GOLD,'gap':RED}.get(c.get('status'),GREY)
        _box(s, col_x, row_y, col_w, Inches(0.82), fill=LGREY)
        _box(s, col_x, row_y, Inches(0.09), Inches(0.82), fill=stc)
        _txt(s, col_x+Inches(0.2), row_y+Inches(0.05), Inches(6), Inches(0.32),
             c.get('name',''), size=12, color=NAVY, bold=True)
        items = c.get('items', [])
        txt = '   '.join(
            (f"{it.get('t','')} {it.get('v','')}" if not it.get('none') else it.get('t',''))
            for it in items[:4])
        _txt(s, col_x+Inches(0.2), row_y+Inches(0.4), Inches(6), Inches(0.38),
             txt, size=10, color=(RED if items and items[0].get('none') else DARK))

    # ---------- P3 AI 진단 ----------
    s = _slide(prs); _header(s, 3, TOTAL, client, 'AI 진단')
    _txt(s, Inches(0.4), Inches(1.25), Inches(6), Inches(0.4),
         '✓ 보유 강점', size=15, color=GREEN, bold=True)
    yy = Inches(1.75)
    for st in rep.get('strength', [])[:5]:
        _txt(s, Inches(0.5), yy, Inches(6), Inches(0.35),
             f"{st.get('h','')} — {st.get('d','')}", size=11, color=DARK)
        yy += Inches(0.42)
    _txt(s, Inches(6.9), Inches(1.25), Inches(6), Inches(0.4),
         '! 보장 공백', size=15, color=RED, bold=True)
    yy = Inches(1.75)
    weak = rep.get('weak', [])
    if not weak:
        _txt(s, Inches(7), yy, Inches(6), Inches(0.35),
             '주요 공백 없음 — 핵심담보 균형 양호', size=11, color=DARK)
    for wk in weak[:5]:
        _txt(s, Inches(7), yy, Inches(6), Inches(0.35),
             f"{wk.get('h','')} — {wk.get('d','')}", size=11, color=DARK)
        yy += Inches(0.42)
    # 보험료 구성 (하단)
    _txt(s, Inches(0.4), Inches(4.2), Inches(6), Inches(0.4),
         '월 보험료 구성', size=13, color=NAVY, bold=True)
    yy = Inches(4.7)
    for b in rep.get('premium_bars', [])[:8]:
        _txt(s, Inches(0.5), yy, Inches(4), Inches(0.3),
             b.get('nm',''), size=10,
             color=(BLUE if b.get('renew') else DARK))
        _txt(s, Inches(4.5), yy, Inches(2), Inches(0.3),
             f"{b.get('amt',0):,}원", size=10, align=PP_ALIGN.RIGHT,
             color=(BLUE if b.get('renew') else DARK))
        yy += Inches(0.3)

    # ---------- P4 CI 선지급 OR 비CI Plan B ----------
    s = _slide(prs); _header(s, 4, TOTAL, client, '핵심 보장 분석')
    if ci_present:
        rate = ci.get('rate', 0)
        _txt(s, Inches(0.4), Inches(1.3), Inches(8), Inches(0.5),
             f'선지급 {rate}% 형', size=20, color=GOLD, bold=True)
        _txt(s, Inches(0.4), Inches(1.95), Inches(12.5), Inches(0.6),
             f"CI 사망보장 {ci.get('samang','')}원. 중대질병 진단 시 {rate}% 선지급, "
             f"잔여 {100-rate}%({ci.get('residual','')}원)는 사망 시 지급된다.",
             size=13, color=DARK)
        # 3개 중대질병 박스
        items = ci.get('items', [])
        bx = Inches(0.5); bw = Inches(3.9); gap = Inches(0.25)
        for i, it in enumerate(items[:3]):
            x = bx + i * (bw + gap)
            _box(s, x, Inches(3.0), bw, Inches(1.6), fill=LGREY, line=GOLD, line_w=Pt(1))
            _txt(s, x, Inches(3.2), bw, Inches(0.4),
                 it.get('t',''), size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, x, Inches(3.7), bw, Inches(0.6),
                 f"{it.get('v','')}", size=22, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, x, Inches(4.35), bw, Inches(0.3),
                 '진단 시 선지급', size=10, color=GREY, align=PP_ALIGN.CENTER)
        _txt(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.2),
             'CI 선지급형 = 중대질병 진단 시 사망보험금 일부를 미리 받아 치료비로 쓰고, '
             '사망 시 잔여분이 지급되는 구조다. 진단 후에도 잔여 사망보장이 유지된다.',
             size=11, color=GREY)
    else:
        # 비CI = Plan B (진단비 정액 지급 구조)
        _txt(s, Inches(0.4), Inches(1.3), Inches(10), Inches(0.5),
             '진단비 정액 지급 구조', size=20, color=NAVY, bold=True)
        _txt(s, Inches(0.4), Inches(1.95), Inches(12.5), Inches(0.7),
             '이 고객은 CI(중대질병 선지급형)가 아닌 일반 진단비 구조다. '
             '암·뇌·심장 진단 시 가입금액이 정액으로 즉시 지급된다(선지급·차감·잔여 없음).',
             size=13, color=DARK)
        # 보유 진단비 3종 박스 (coverage에서 암·뇌·심)
        picks = []
        for c in rep.get('coverage', []):
            nm = c.get('name','')
            if any(k in nm for k in ('암','뇌혈관','심장')):
                its = [it for it in c.get('items',[]) if not it.get('none')]
                if its:
                    picks.append((nm.split(' ')[0].replace('(＋빈맥)',''), its[0]))
        bx = Inches(0.5); bw = Inches(3.9); gap = Inches(0.25)
        for i, (nm, it) in enumerate(picks[:3]):
            x = bx + i * (bw + gap)
            _box(s, x, Inches(3.0), bw, Inches(1.6), fill=LGREY, line=NAVY, line_w=Pt(1))
            _txt(s, x, Inches(3.2), bw, Inches(0.4),
                 nm, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, x, Inches(3.7), bw, Inches(0.6),
                 f"{it.get('v','')}", size=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, x, Inches(4.35), bw, Inches(0.3),
                 '진단 즉시 정액 100%', size=10, color=GREY, align=PP_ALIGN.CENTER)
        _txt(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.0),
             '진단비는 실제 치료비와 무관하게 약정된 가입금액이 정액 지급된다. '
             '치료가 장기화되면 주요치료비(정액형)로 보완이 필요하다.',
             size=11, color=GREY)

    # ---------- P5 부위별 충족률 ----------
    s = _slide(prs); _header(s, 5, TOTAL, client, '부위별 충족률')
    donuts = rep.get('donut_detail', [])
    cols = 5; dw = Inches(2.4); dh = Inches(1.4)
    x0 = Inches(0.5); y0 = Inches(1.4)
    for i, d in enumerate(donuts[:10]):
        r = i // cols; c = i % cols
        x = x0 + c * dw; yy = y0 + r * Inches(1.7)
        pc = _pct_color(d.get('pct',0))
        _box(s, x, yy, Inches(2.1), Inches(1.1), fill=LGREY, line=pc, line_w=Pt(2))
        _txt(s, x, yy+Inches(0.12), Inches(2.1), Inches(0.5),
             f"{d.get('pct',0)}%", size=22, color=pc, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, x, yy+Inches(0.68), Inches(2.1), Inches(0.35),
             d.get('name',''), size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # 산정 근거
    _txt(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4),
         '충족률 = 보유 ÷ 연령밴드 권장액 × 100 (상한 100%). '
         f"{rep.get('band_label','40대')} 표준밴드 적용. 녹색 충실(70%↑) / 골드 보강(40~69%) / 빨강 취약(40%↓).",
         size=10, color=GREY)

    # ---------- P6 담보별 보장범위 (뇌·심 개별축·산정특례) ----------
    s = _slide(prs); _header(s, 6, TOTAL, client, '담보별 보장범위')
    _txt(s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.4),
         '뇌·심 담보축 = 각각 개별 담보·각각 보상 · 산정특례·순환계·외상성뇌출혈 = 단독',
         size=12, color=NAVY, bold=True)
    # 뇌
    _txt(s, Inches(0.4), Inches(1.75), Inches(6), Inches(0.35),
         '뇌 — 뇌혈관진단비 · 순환계 · 산정특례 (각 개별)', size=12, color=NAVY, bold=True)
    _txt(s, Inches(0.5), Inches(2.2), Inches(6.2), Inches(2.5),
         '· 출혈성(I60~62)·허혈성(I63~66)·기타(I64·67~69)\n'
         '· 산정특례 = 뇌혈관질환 I60~69 + Q28 + S06 전체\n'
         '· 외상성뇌출혈(S06) = 뇌혈관진단비 미보장, 산정특례 축only\n'
         '· 각 축 각각 개별 담보로 각각 보상', size=11, color=DARK)
    # 심
    _txt(s, Inches(6.9), Inches(1.75), Inches(6), Inches(0.35),
         '심장 — 허혈성 · 심장특정 · 순환계 · 산정특례 (각 개별)', size=12, color=NAVY, bold=True)
    _txt(s, Inches(7.0), Inches(2.2), Inches(6.0), Inches(2.5),
         '· 허혈성(I20~25)·판막·염증·부정맥·심부전·심근병증\n'
         '· 산정특례 = 심혈관질환 I20~50 + 판막 전체\n'
         '· 빈맥(I47·48) = 마스터 무행·전 묶음 제외(고정)\n'
         '· 각 축 각각 개별 담보로 각각 보상', size=11, color=DARK)
    # 산정특례 기준박스
    _box(s, Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.7), fill=RGBColor(0xFD,0xF9,0xEF), line=GOLD, line_w=Pt(1))
    _txt(s, Inches(0.6), Inches(5.15), Inches(12), Inches(0.4),
         '산정특례 기준 (진단 기반 · 별개 담보축)', size=12, color=RGBColor(0x8A,0x5D,0x00), bold=True)
    _txt(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.0),
         '산정특례 = 위 범위와 축이 다른 별개 담보 — 마스터 산정특례심장·산정특례뇌혈관 전용행에서 진단코드 기반 지급. '
         '대상 코드범위 [뇌] I60~69·Q28·S06 / [심] I20~50·판막. 각각 개별 담보로 각각 보상. '
         '지급조건·기간(30일·5% 등)은 회사·약관별 [확인].', size=10, color=RGBColor(0x6A,0x5A,0x20))

    # ---------- P7 주요치료비 정리 ----------
    s = _slide(prs); _header(s, 7, TOTAL, client, '주요치료비 변천사')
    _txt(s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.4),
         '3대 주요치료비 담보의 변천사  ·  병원비형 → 정액형 → 확대형 → 생활비 · 정액형 가입금액 100만원부터', size=13, color=NAVY, bold=True)
    _box(s, Inches(0.4), Inches(1.75), Inches(6.25), Inches(0.42), fill=RGBColor(0xB9,0x54,0x0B))
    _txt(s, Inches(0.55), Inches(1.79), Inches(6), Inches(0.35), '\u25a0 \uc554 \uc8fc\uc694\uce58\ub8cc\ube44  \u00b7  4\uc138\ub300 \uc9c4\ud654', size=12, color=WHITE, bold=True)
    _txt(s, Inches(0.5), Inches(2.3), Inches(6.1), Inches(0.7),
         '\u2460\uad6c\uac04\ud615(~24.11 \ub2e8\uc885) \u2192 \u2461\uc815\uc561\ud615(24.11~, 100\ub9cc~) \u2192 \u2462\ud558\uc774\ud074\ub798\uc2a4(25.03~ \ube44\uae09\uc5ec \uc804\uc561\ubcf8\uc778) \u2192 \u2463\uc0dd\ud65c\ube44(25.08~)', size=10, color=DARK)
    _box(s, Inches(0.5), Inches(3.0), Inches(6.05), Inches(1.15), fill=LGREY, line=GREY, line_w=Pt(0.75))
    _txt(s, Inches(0.62), Inches(3.08), Inches(5.9), Inches(1.05),
         '\u00b7 \uad6c\uac04\ud615 12\ub2e8\uacc4: \uce58\ub8cc\ube44 300\ub9cc~1\uc5b5 \uad6c\uac04 \ud558\ud55c \ub3c4\ub2ec \uc2dc \uadf8 \uad6c\uac04 \uc815\uc561. \uc2e4\uc190 \uc911\ubcf5 \uac00\ub2a5 \u00b7 2024.11 \ub2e8\uc885\n\u00b7 \uc815\uc561\ud615: \uce58\ub8cc \uc0ac\uc2e4\ub9cc\uc73c\ub85c \uc815\uc561(\uba74\ucc45 90\uc77c) \u00b7 \ud558\uc774\ud074\ub798\uc2a4: \uc591\uc131\uc790\u00b7\uc911\uc785\uc790\u00b7\ud45c\uc801\u00b7\uba74\uc5ed \ucee4\ubc84(2\ucc9c\ub9cc\u00d710\ub144)', size=9.5, color=DARK)
    _box(s, Inches(0.5), Inches(4.3), Inches(6.05), Inches(0.75), fill=RGBColor(0xFB,0xEC,0xEC), line=RED, line_w=Pt(0.75))
    _txt(s, Inches(0.62), Inches(4.38), Inches(5.9), Inches(0.65),
         '\u25b2 5\uc138\ub300 \uc2e4\uc190 \ube44\uae09\uc5ec = \uc790\uae30\ubd80\ub2f4 50%\u00b7\ud55c\ub3c4 1,000\ub9cc(\ucd95\uc18c). \ud45c\uc801\u00b7\uba74\uc5ed \uace0\uac00\uc5d4 \ubd80\uc871 \u2192 \uc8fc\uc694\uce58\ub8cc\ube44\u00b7\ud558\uc774\ud074\ub798\uc2a4\uac00 \uba54\uc6c0', size=9.5, color=RED)
    _box(s, Inches(0.5), Inches(5.2), Inches(6.05), Inches(1.35), fill=RGBColor(0xEA,0xF2,0xFB), line=RGBColor(0x14,0x56,0xB0), line_w=Pt(0.75))
    _txt(s, Inches(0.62), Inches(5.28), Inches(5.9), Inches(1.25),
         '\u25cf \uc0c1\ub2f4: "\ud45c\uc801\ud56d\uc554 \ud55c \ub2ec \uc218\ubc31\ub9cc \uc6d0\uc774 \uc804\uc561 \ubcf8\uc778\ubd80\ub2f4\uc774\uc5d0\uc694. 5\uc138\ub300 \uc2e4\uc190\uc740 \ube44\uae09\uc5ec\uac00 1\ucc9c\ub9cc\uc73c\ub85c \uc904\uc5c8\uc2b5\ub2c8\ub2e4. \uad6c\uac04\ud615\uc740 \uc774\uc81c \ubabb \ub4dc\ub294 \ub2f4\ubcf4\ub77c \uc720\uc9c0\uac00 \ub2f5\uc774\uace0, \uc5c6\uc73c\uba74 \uc815\uc561\ud615\u00b7\ud558\uc774\ud074\ub798\uc2a4\ub85c \ucc44\uc6cc\uc57c \ud569\ub2c8\ub2e4."', size=9.5, color=NAVY)
    _box(s, Inches(6.9), Inches(1.75), Inches(6.05), Inches(0.42), fill=RGBColor(0x1F,0x5F,0xA8))
    _txt(s, Inches(7.05), Inches(1.79), Inches(6), Inches(0.35), '\u25a0 \ub1cc\u00b7\uc2ec\uc7a5 \uc8fc\uc694\uce58\ub8cc\ube44  \u00b7  \uc554\uacfc \ub3d9\uc77c \uc9c4\ud654', size=12, color=WHITE, bold=True)
    _txt(s, Inches(7.0), Inches(2.3), Inches(5.9), Inches(0.7),
         '\u2460\ubcd1\uc6d0\ube44\ud615(~24.11 \ub2e8\uc885) \u2192 \u2461\uc815\uc561\ud615 2\ub300(24.11~ \uba54\ub9ac\uce20 \ucd5c\ucd08) \u2192 \u2462\uc21c\ud658\uacc4(25.01~ \ubd80\uc815\ub9e5\u00b7\uc2ec\ubd80\uc804\u00b7\ub3d9\ub9e5\ub958) \u2192 \u2463\uc21c\ud658\uacc4 \uc0dd\ud65c\ube44(25.06~)', size=10, color=DARK)
    _box(s, Inches(7.0), Inches(3.0), Inches(5.9), Inches(1.15), fill=LGREY, line=GREY, line_w=Pt(0.75))
    _txt(s, Inches(7.12), Inches(3.08), Inches(5.75), Inches(1.05),
         '\u00b7 \ud68c\uc0ac\ubcc4 \ubc94\uc704: \ud3ec\uad04\ud615(\uc0bc\uc131\u00b7\ud604\ub300\u00b7\uba54\ub9ac\uce20\u00b7\ud55c\ud654\u00b7\ub86f\ub370\u00b7\ub3d9\uc591) / \uc21c\ud658\uacc4 \ub113\uc74c(DB\u00b7KB\u00b7\ud765\uad6d) / \ud2b9\uc815\uc21c\ud658\uacc4(NH\ub18d\ud611)\n\u00b7 \ub1cc\u00b7\uc2ec \uad6c\uac04\ud615 = \uae09\uc5ec \ubcf8\uc778\ubd80\ub2f4 \uae30\uc900(100\ub9cc \ubbf8\ub9cc \uba74\ucc45 / 3,000\ub9cc \ud55c\ub3c4) \u00b7 2024.11 \ub2e8\uc885', size=9.5, color=DARK)
    _box(s, Inches(7.0), Inches(4.3), Inches(5.9), Inches(0.75), fill=RGBColor(0xFD,0xEC,0xEC), line=RED, line_w=Pt(0.75))
    _txt(s, Inches(7.12), Inches(4.38), Inches(5.75), Inches(0.65),
         '\u25b2 \uc7ac\ubc1c\ub960: \ub1cc\uacbd\uc0c9(\ud5c8\ud608\uc131) 5\ub144 \ub0b4 20~30% \uc7ac\ubc1c \u00b7 \uc2ec\uc7a5 \uc2a4\ud150\ud2b8 5\ub144 \uc7ac\ud611\ucc29 15% \u2192 \uc8fc\uc694\uce58\ub8cc\ube44\ub294 \uc7ac\ubc1c\u00b7\uc7a5\uae30\uce58\ub8cc \ubc18\ubcf5 \ubcf4\uc7a5', size=9.5, color=RED)
    _box(s, Inches(7.0), Inches(5.2), Inches(5.9), Inches(1.35), fill=RGBColor(0xEA,0xF2,0xFB), line=RGBColor(0x14,0x56,0xB0), line_w=Pt(0.75))
    _txt(s, Inches(7.12), Inches(5.28), Inches(5.75), Inches(1.25),
         '\u25cf \uc0c1\ub2f4: "\ub1cc\u00b7\uc2ec\uc7a5\uc740 30\uc77c\uc740 \uad6d\uac00\uac00 \ub3d5\uc9c0\ub9cc, \uc2a4\ud150\ud2b8 \ud6c4 1\ub144 \uc57d\ubb3c, \ub1cc\uc878\uc911 \uc7ac\ud65c\uc740 \uba87 \ub144\uc529 \uac11\ub2c8\ub2e4. \uac8c\ub2e4\uac00 5\ub144 \ub0b4 2~3\uba85\uc740 \uc7ac\ubc1c\ud574\uc694. \uadf8 \uae34 \uce58\ub8cc\ube44\ub97c \uc8fc\uc694\uce58\ub8cc\ube44\uac00 \ubc18\ubcf5\ud574 \ucc44\uc6c1\ub2c8\ub2e4."', size=9.5, color=NAVY)
    _txt(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
         '\u203b \uc138\ub300 \ud310\ubcc4 = \uacc4\uc57d\uc77c \uae30\uc900 \u00b7 \uad6c\uac04\ud615(~24.11) \ubcf4\uc720\uc790 = \ub2e8\uc885 \ub2f4\ubcf4, \ud574\uc9c0 \uae08\uc9c0 \u00b7 \uadfc\uac70 \uad50\uc721\uc790\ub8cc 2607 + \ubcf4\ud5d8\uc800\ub110\u00b7\uae08\uc735\uc704(2026.07) \u00b7 \uc9c4\ub2e8\ube44 \ucd95\uacfc \ubcc4\uac1c', size=8.5, color=GREY)


    # ---------- P8 상담 워크시트 (산정특례 칸 포함) ----------
    s = _slide(prs); _header(s, 8, TOTAL, client, '상담 워크시트')
    _txt(s, Inches(0.4), Inches(1.2), Inches(8), Inches(0.4),
         '지금 고객의 3대 주요치료비는?', size=16, color=NAVY, bold=True)
    # 암 / 뇌심 2열
    _txt(s, Inches(0.4), Inches(1.75), Inches(6), Inches(0.35),
         '암', size=13, color=GOLD, bold=True)
    _txt(s, Inches(0.5), Inches(2.2), Inches(6), Inches(1.5),
         '✔ 암 진단비 — 걸렸을 때 일시금 (기본)\n'
         '✘ 암 주요치료비 — 수술·방사선·약물 정액 (100만~)\n'
         '□ 하이클래스(비급여) — 표적·면역·중입자\n'
         '□ 암 생활비 — 치료 중 소득보상', size=11, color=DARK)
    _txt(s, Inches(6.9), Inches(1.75), Inches(6), Inches(0.35),
         '뇌 · 심장', size=13, color=GOLD, bold=True)
    _txt(s, Inches(7.0), Inches(2.2), Inches(6), Inches(1.5),
         '✔ 뇌·심 진단비 — 뇌혈관·허혈성 일시금\n'
         '✘ 2대 주요치료비 — 수술·혈전용해·중환자실\n'
         '□ 순환계 주요치료비 — 부정맥·심부전 확대\n'
         '□ 순환계 생활비 — 치료 중 소득보상', size=11, color=DARK)
    # 산정특례 별개축 박스
    _box(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(1.5), fill=RGBColor(0xFD,0xF9,0xEF), line=GOLD, line_w=Pt(1.5))
    _txt(s, Inches(0.6), Inches(4.15), Inches(12), Inches(0.35),
         '산정특례 — 뇌·심 각각 개별 담보 · 진단만으로 지급', size=12, color=RGBColor(0x8A,0x5D,0x00), bold=True)
    _txt(s, Inches(0.6), Inches(4.6), Inches(6), Inches(0.8),
         '산정특례(뇌) — 뇌혈관질환 I60~69 · Q28 · S06\n상태: ____________  만원', size=11, color=RGBColor(0x6A,0x5A,0x20))
    _txt(s, Inches(6.9), Inches(4.6), Inches(6), Inches(0.8),
         '산정특례(심장) — 심혈관질환 I20~50 · 판막 전체\n상태: ____________  만원', size=11, color=RGBColor(0x6A,0x5A,0x20))
    _txt(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0),
         '진단비(일시금)와 주요치료비(치료 실비·정액)는 별개 축이라 둘 다 필요하다. '
         '산정특례는 진단만으로 지급되는 개별 담보다.', size=10, color=GREY)

    prs.save(out)
    return out


if __name__ == '__main__':
    import sys
    from coverage_benchmark import map_excel_to_report
    xl = sys.argv[1] if len(sys.argv) > 1 else '보장진단_정기철.xlsx'
    rep = map_excel_to_report(xl, settings={'client': '테스트', 'branch': '온빛센터 바름지점',
                                            'manager': '최은혜', 'title': '지점장', 'phone': ''})
    build_report_pptx(rep, 'test_report.pptx')
    print('생성 완료: test_report.pptx')
