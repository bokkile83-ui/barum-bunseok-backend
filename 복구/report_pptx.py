# -*- coding: utf-8 -*-
"""BARUM 보장진단서 PPT v33 — 도형·텍스트 기반(편집 가능).
   지점장 지시(2026.07.08): PDF→슬라이드 이미지 방식 폐기. 값 텍스트를 직접 수정할 수 있어야 한다.
   · 값 출처 = rep 딕셔너리(= coverage_benchmark 가 완성 엑셀에서 산출). 등식2 유지.
   · 게이지 링만 픽셀(PIL). 그 외 전부 도형·텍스트.
   · 폰트 = 맑은 고딕 (latin + ea 동시 지정).
"""
import os, math, tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
GOLD  = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY  = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT = RGBColor(0xF2, 0xF3, 0xF5)
LINE  = RGBColor(0xDD, 0xDD, 0xDD)
GREEN = RGBColor(0x1D, 0x7A, 0x46)
RED   = RGBColor(0xC0, 0x39, 0x2B)
BLUE  = RGBColor(0x18, 0x5F, 0xA5)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
FONT  = "맑은 고딕"

W = Emu(int(297 / 25.4 * 914400))
H = Emu(int(210 / 25.4 * 914400))


def _setfont(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)


def box(s, x, y, w, h, text, size=10, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    sh = s.shapes.add_textbox(x, y, w, h)
    tf = sh.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.margin_top = tf.margin_bottom = Emu(18288)
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        _setfont(r)
    return sh


def rect(s, x, y, w, h, fill, line=None, rad=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    else:    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def header(s, name, page, sub):
    rect(s, 0, 0, W, Inches(1.0), NAVY)
    rect(s, 0, Inches(1.0), W, Inches(0.03), GOLD)
    box(s, Inches(0.42), Inches(0.10), Inches(6), Inches(0.26), "MAKEONE · 보장분석 리포트", 8.5, True, GOLD)
    box(s, Inches(0.42), Inches(0.36), Inches(7.5), Inches(0.48), f"{name} 고객님 보장 진단서", 18, True, WHITE)
    box(s, W - Inches(1.6), Inches(0.14), Inches(1.2), Inches(0.44), str(page), 22, True, WHITE, PP_ALIGN.RIGHT)
    box(s, W - Inches(3.2), Inches(0.62), Inches(2.8), Inches(0.26), sub, 8, False, GOLD, PP_ALIGN.RIGHT)


def footer(s, name, page, total):
    rect(s, 0, H - Inches(0.38), W, Inches(0.38), NAVY)
    box(s, Inches(0.42), H - Inches(0.34), Inches(4), Inches(0.26), "MAKEONE 보장분석 자동화", 7.5, True, GOLD)
    box(s, W - Inches(4.4), H - Inches(0.34), Inches(4), Inches(0.26),
        f"{name} 고객님 · {page} / {total}", 7.5, False, WHITE, PP_ALIGN.RIGHT)


def _gauge_png(pct, path, px=360):
    from PIL import Image, ImageDraw
    S = px * 4
    img = Image.new("RGBA", (S, S), (0, 0, 0, 0)); d = ImageDraw.Draw(img)
    col = (0x1D, 0x7A, 0x46) if pct >= 70 else (0xC9, 0xA2, 0x27) if pct >= 40 else (0xC0, 0x39, 0x2B)
    ring = int(S * 0.085); pad = ring // 2 + 4
    bb = [pad, pad, S - pad, S - pad]
    d.arc(bb, 0, 360, fill=(0xE3, 0xE5, 0xE9, 255), width=ring)
    if pct > 0:
        d.arc(bb, -90, -90 + 360 * pct / 100, fill=col + (255,), width=ring)
        if pct < 100:
            a = math.radians(-90 + 360 * pct / 100); r = (S - 2 * pad) / 2
            cx, cy = S / 2 + r * math.cos(a), S / 2 + r * math.sin(a)
            d.ellipse([cx - ring / 2, cy - ring / 2, cx + ring / 2, cy + ring / 2], fill=col + (255,))
    img.resize((px, px), Image.LANCZOS).save(path)
    return path


def table(s, x, y, cols, rows, widths, hh=0.26, rh=0.24, fs=8):
    cx = x
    for j, cw in enumerate(widths):
        rect(s, cx, y, Inches(cw), Inches(hh), NAVY)
        box(s, cx, y, Inches(cw), Inches(hh), cols[j], fs, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, False)
        cx += Inches(cw)
    yy = y + Inches(hh)
    for i, row in enumerate(rows):
        cx = x; bg = WHITE if i % 2 == 0 else LIGHT
        for j, cw in enumerate(widths):
            cell = row[j]
            txt, col = cell if isinstance(cell, tuple) else (cell, BLACK)
            rect(s, cx, yy, Inches(cw), Inches(rh), bg, LINE)
            box(s, cx, yy, Inches(cw), Inches(rh), txt, fs, False, col,
                PP_ALIGN.CENTER if j else PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, False)
            cx += Inches(cw)
        yy += Inches(rh)
    return yy


def build_report_pptx(rep, out, dpi=150):
    """rep(=coverage_benchmark 산출) → 도형·텍스트 10슬라이드. dpi 인자는 호환용(무시)."""
    name = rep.get('client', '고객')
    TOT = 9
    tmp = tempfile.mkdtemp()
    prs = Presentation(); prs.slide_width, prs.slide_height = W, H
    BL = prs.slide_layouts[6]

    # 1 표지
    s = prs.slides.add_slide(BL)
    rect(s, 0, 0, W, H, NAVY)
    box(s, Inches(0.9), Inches(2.4), Inches(6), Inches(0.34), "M A K E O N E", 13, True, GOLD)
    box(s, Inches(0.9), Inches(2.8), Inches(8), Inches(0.9), "보장 진단서", 36, True, WHITE)
    box(s, Inches(0.9), Inches(3.8), Inches(8), Inches(0.5), f"{name} 고객님", 16, False, WHITE)
    rect(s, Inches(0.95), Inches(4.6), Inches(2.6), Inches(0.02), GOLD)
    for i, (k, v) in enumerate([("보유 계약", f"{rep.get('n_contract',0)} 건"),
                                ("월 납입보험료", f"{rep.get('premium',0):,} 원"),
                                ("갱신 / 비갱신", f"{rep.get('renew',0)} / {rep.get('nonrenew',0)}")]):
        x = Inches(0.9 + i * 3.3)
        box(s, x, Inches(4.9), Inches(3.0), Inches(0.26), k, 8.5, False, GOLD)
        box(s, x, Inches(5.16), Inches(3.0), Inches(0.42), v, 15, True, WHITE)
    box(s, Inches(0.9), H - Inches(0.9), Inches(8), Inches(0.28),
        "MAKEONE · 보장분석 자동화 리포트", 8, False, GREY)

    # 2 보장 현황
    s = prs.slides.add_slide(BL); header(s, name, 1, "보장 현황")
    for i, c in enumerate(rep.get('coverage', [])[:10]):
        col, row = i % 2, i // 2
        x = Inches(0.42 + col * 5.62); y = Inches(1.28 + row * 0.86)
        rect(s, x, y, Inches(1.35), Inches(0.74), NAVY)
        box(s, x, y, Inches(1.35), Inches(0.74), c['name'], 8.5, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        rect(s, x + Inches(1.35), y, Inches(4.1), Inches(0.74), LIGHT, LINE)
        it = [f"{i['t']} {i['v']}" for i in c.get('items', [])][:6]
        txt = "   ".join(it[:2]) + (("\n" + "   ".join(it[2:4])) if len(it) > 2 else "")
        box(s, x + Inches(1.4), y, Inches(4.0), Inches(0.74),
            txt if it else "미가입", 8, False, BLACK if it else RED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, name, 1, TOT)

    # 3 AI 진단
    s = prs.slides.add_slide(BL); header(s, name, 2, "AI 진단")
    box(s, Inches(0.42), Inches(1.22), Inches(5), Inches(0.3), "갱신형 — 보험료 인상 가능", 11, True, BLUE)
    box(s, Inches(5.9), Inches(1.22), Inches(5), Inches(0.3), "비갱신형 — 만기까지 고정", 11, True, RED)
    for lst, x0, col in ((rep.get('renew_list', []), 0.42, BLUE), (rep.get('nonrenew_list', []), 5.9, BLACK)):
        for j, c in enumerate(lst[:7]):
            y = Inches(1.6 + j * 0.5); x = Inches(x0)
            rect(s, x, y, Inches(5.15), Inches(0.44), LIGHT, LINE)
            box(s, x + Inches(0.08), y, Inches(3.6), Inches(0.44), c['nm'][:30], 8, False, BLACK, anchor=MSO_ANCHOR.MIDDLE)
            box(s, x + Inches(3.7), y, Inches(1.35), Inches(0.44), c['v'], 9, True, col, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
    y0 = Inches(1.6 + 7 * 0.5 + 0.15)
    box(s, Inches(0.42), y0, Inches(5), Inches(0.28), "보유 강점", 10, True, GREEN)
    for j, a in enumerate(rep.get('strength', [])[:5]):
        box(s, Inches(0.5), y0 + Inches(0.3 + j * 0.24), Inches(5), Inches(0.22), f"· {a['h']} — {a['d']}", 8, False, BLACK)
    box(s, Inches(5.9), y0, Inches(5), Inches(0.28), "보장 공백", 10, True, RED)
    wk = rep.get('weak', [])
    if wk:
        for j, a in enumerate(wk[:5]):
            box(s, Inches(5.98), y0 + Inches(0.3 + j * 0.24), Inches(5), Inches(0.22), f"· {a['h']} — {a['d']}", 8, False, BLACK)
    else:
        box(s, Inches(5.98), y0 + Inches(0.3), Inches(5), Inches(0.22), "· 주요 공백 없음 — 핵심담보 균형이 양호합니다.", 8, False, GREY)
    footer(s, name, 2, TOT)

    # 4 핵심 보장 분석 (CI 3-state)
    s = prs.slides.add_slide(BL); header(s, name, 3, "핵심 보장 분석")
    ci = rep.get('ci', {}); st = ci.get('status') or ('ci' if ci.get('present') else 'none')
    if st == 'ci':
        rate = ci.get('rate', 0); sam = ci.get('samang', ''); res = ci.get('residual', '')
        box(s, Inches(0.42), Inches(1.2), Inches(6), Inches(0.32), "CI 선지급 분석  CRITICAL ILLNESS · PRE-PAYMENT", 11, True, NAVY)
        rect(s, Inches(0.42), Inches(1.56), W - Inches(0.84), Inches(2.35), WHITE, GOLD)
        rect(s, Inches(0.62), Inches(1.72), Inches(1.7), Inches(0.32), NAVY, rad=True)
        box(s, Inches(0.62), Inches(1.72), Inches(1.7), Inches(0.32), f"선지급 {rate}% 형", 9, True, GOLD, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, False)
        box(s, Inches(0.62), Inches(2.1), W - Inches(1.3), Inches(0.34),
            f"CI 주계약 사망보험금 {sam}. 중대질병 진단 시 {rate}% 선지급, 잔여 {100-rate}%({res})는 사망 시 지급된다.", 9, False, BLACK)
        rect(s, Inches(0.62), Inches(2.5), Inches(4.9), Inches(0.62), NAVY)
        box(s, Inches(0.75), Inches(2.56), Inches(4.6), Inches(0.22), f"선지급 {rate}%", 8, True, GOLD)
        box(s, Inches(0.75), Inches(2.76), Inches(4.6), Inches(0.3), f"{ci['items'][0]['v'] if ci.get('items') else ''} · 진단 시", 12, True, WHITE)
        rect(s, Inches(5.72), Inches(2.5), Inches(4.9), Inches(0.62), LIGHT, LINE)
        box(s, Inches(5.85), Inches(2.56), Inches(4.6), Inches(0.22), f"잔여 {100-rate}%", 8, True, GREY)
        box(s, Inches(5.85), Inches(2.76), Inches(4.6), Inches(0.3), f"{res} · 사망 시", 12, True, NAVY)
        for j, it in enumerate(ci.get('items', [])[:3]):
            x = Inches(0.62 + j * 3.42)
            rect(s, x, Inches(3.2), Inches(3.2), Inches(0.6), WHITE, GOLD)
            box(s, x, Inches(3.24), Inches(3.2), Inches(0.22), it['t'], 8, True, GREY, PP_ALIGN.CENTER)
            box(s, x, Inches(3.44), Inches(3.2), Inches(0.32), it['v'], 12, True, GOLD, PP_ALIGN.CENTER)
        box(s, Inches(0.62), Inches(3.86), W - Inches(1.3), Inches(0.3),
            "■ CI 선지급형 — 중대질병 진단 시 사망보험금 일부를 미리 받아 치료비로 쓰고, 사망 시 잔여분이 지급된다. 선지급률은 50% 또는 80%만 존재한다.",
            7.5, False, GREY)
    elif st == 'check':
        box(s, Inches(0.42), Inches(1.2), Inches(6), Inches(0.32), "핵심 보장 분석  CI 판정 보류 · CONFIRM", 11, True, NAVY)
        rect(s, Inches(0.42), Inches(1.56), W - Inches(0.84), Inches(1.5), RGBColor(0xFB, 0xF6, 0xE6), GOLD)
        box(s, Inches(0.65), Inches(1.7), Inches(9), Inches(0.36), "CI 판정 [확인]", 15, True, GOLD)
        box(s, Inches(0.65), Inches(2.12), W - Inches(1.4), Inches(0.8),
            "■ 판정 보류 — 수기 확인 필요\n상품명에 CI · GI · 리빙케어가 있으나 중대한OO 담보 금액이 추출되지 않았다.\nCI 여부와 선지급률(50%/80%)을 약관으로 직접 확인한 뒤 확정하라. 추측 금지 · [확인].", 9, False, BLACK)
    else:
        box(s, Inches(0.42), Inches(1.2), Inches(6), Inches(0.32), "핵심 보장 분석  GENERAL INSURANCE · NON-CI", 11, True, NAVY)
        rect(s, Inches(0.42), Inches(1.56), W - Inches(0.84), Inches(1.3), LIGHT, LINE)
        box(s, Inches(0.65), Inches(1.7), Inches(9), Inches(0.36), "CI보험 아님", 15, True, BLUE)
        box(s, Inches(0.65), Inches(2.1), W - Inches(1.4), Inches(0.66),
            "■ 일반 보험기준 — 이 고객은 CI(중대질병 선지급형) 계약이 아니다. 암·뇌·심 진단비는 진단 즉시 정액 100% "
            "지급되는 일반형이며, 사망보험금 차감·선지급·잔여 개념이 없다.", 9, False, BLACK)

    yb = Inches(4.35) if st == 'ci' else Inches(3.1)
    box(s, Inches(0.42), yb, Inches(6), Inches(0.3), "주요 치료비 정리  TREATMENT BENEFITS", 10.5, True, NAVY)
    for j, c in enumerate(rep.get('chiryo', [])[:5]):
        x = Inches(0.42 + j * 2.14); ok = c['value'] != '미가입'
        rect(s, x, yb + Inches(0.34), Inches(1.98), Inches(0.72), WHITE, GOLD if ok else LINE)
        box(s, x, yb + Inches(0.4), Inches(1.98), Inches(0.24), c['name'], 7.5, True, GREY, PP_ALIGN.CENTER)
        box(s, x, yb + Inches(0.64), Inches(1.98), Inches(0.32), c['value'], 11 if ok else 9, True, NAVY if ok else RED, PP_ALIGN.CENTER)
    yc = yb + Inches(1.2)
    box(s, Inches(0.42), yc, Inches(6), Inches(0.3), "보장 진단 코멘트  SUMMARY", 10.5, True, NAVY)
    rect(s, Inches(0.42), yc + Inches(0.32), W - Inches(0.84), Inches(0.9), LIGHT, LINE)
    cmt = f"{name} 고객님은 보유 {rep.get('n_contract',0)}건 · 월 {rep.get('premium',0):,}원의 보장을 운용하고 있습니다."
    if st == 'ci':
        cmt += f" CI 선지급형 보유로 중대질병 진단 시 진단자금이 즉시 지급되며, 진단 후에도 잔여 사망보장 {ci.get('residual','')}이 유지됩니다."
    cmt += f" 보강이 필요한 공백 영역은 {rep.get('gap_count',0)}개입니다."
    box(s, Inches(0.62), yc + Inches(0.4), W - Inches(1.3), Inches(0.7), cmt, 9, False, BLACK)
    footer(s, name, 3, TOT)

    # 5 부위별 충족률
    s = prs.slides.add_slide(BL); header(s, name, 4, "부위별 충족률")
    for i, g in enumerate(rep.get('donuts', [])[:10]):
        cx = Inches(1.15 + (i % 5) * 2.15); cy = Inches(1.85 + (i // 5) * 1.6); d = Inches(1.0)
        p = _gauge_png(g['pct'], os.path.join(tmp, f"g{i}.png"))
        s.shapes.add_picture(p, cx - d // 2, cy - d // 2, d, d)
        box(s, cx - d // 2, cy - Inches(0.14), d, Inches(0.28), f"{g['pct']}%", 11, True, BLACK, PP_ALIGN.CENTER, wrap=False)
        box(s, cx - Inches(0.8), cy + d // 2 + Inches(0.02), Inches(1.6), Inches(0.24), g['name'], 8, True, BLACK, PP_ALIGN.CENTER, wrap=False)
    box(s, Inches(0.42), Inches(5.0), Inches(6), Inches(0.28),
        f"충족률 산정 근거   보유 ÷ {rep.get('band_label','40대')} 권장", 10, True, NAVY)
    dd = rep.get('donut_detail', [])
    def rows(sub):
        return [(d['name'], d['have'], d['rec'],
                 (f"{d['pct']}%", GREEN if d['pct'] >= 70 else (GOLD if d['pct'] >= 40 else RED))) for d in sub]
    table(s, Inches(0.42), Inches(5.32), ["영역", "보유", "권장", "충족률"], rows(dd[:5]), [1.6, 1.3, 1.3, 0.95], rh=0.2, fs=7.5)
    table(s, Inches(5.9), Inches(5.32), ["영역", "보유", "권장", "충족률"], rows(dd[5:]), [1.6, 1.3, 1.3, 0.95], rh=0.2, fs=7.5)
    for j, a in enumerate(rep.get('advice', [])[:1]):
        box(s, Inches(0.42), Inches(6.6), W - Inches(0.84), Inches(0.5), f"■ {a['t']}\n{a['d']}", 7.5, False, GREY)
    footer(s, name, 4, TOT)

    # 6 담보별 보장범위
    s = prs.slides.add_slide(BL); header(s, name, 5, "담보별 보장범위")
    sb = set(rep.get('scope_brain', [])); sh_ = set(rep.get('scope_heart', []))
    O, X = "●", "○"
    box(s, Inches(0.42), Inches(1.2), Inches(5), Inches(0.28), "뇌 — 질병코드별 커버", 10.5, True, NAVY)
    brain = [("뇌출혈  I60~62", O if 'hem' in sb else X, O, O),
             ("뇌졸증·뇌경색  I63·65·66", O if 'infarct' in sb else X, O, O),
             ("기타 뇌혈관질환  I64·67~69", O if 'other' in sb else X, O, O),
             ("뇌동맥류·정맥류  I71·72", X, O, X),
             ("선천 뇌혈관기형  Q28.0~28.3", X, X, O),
             ("외상성 뇌출혈  S06", X, X, O)]
    table(s, Inches(0.42), Inches(1.52), ["질병 (코드)", "뇌혈관진단비", "순환계", "산정특례"], brain, [2.5, 1.2, 0.85, 0.85], rh=0.25, fs=8)
    box(s, Inches(0.42), Inches(3.4), Inches(5.3), Inches(0.6),
        "· 외상성 뇌출혈(S06) = 뇌혈관진단비 미보장 → 산정특례 축only\n· 산정특례 = 진단 기반 별개 담보축 · 지급조건 [확인]", 7.5, False, GREY)
    box(s, Inches(5.9), Inches(1.2), Inches(5), Inches(0.28), "심장 — 질병코드별 커버", 10.5, True, NAVY)
    heart = [("급성심근경색  I21~23", O if 'ami' in sh_ else X, O, O, O),
             ("협심증  I20", O if 'angina' in sh_ else X, O, O, O),
             ("기타·만성 허혈  I24·25", O if 'chronic' in sh_ else X, O, O, O),
             ("심장판막  I05·I34~37", X, O, O, O),
             ("심근·심내막 염증  I30~33·I40", X, O, O, O),
             ("부정맥  I49", O if 'arrhy' in sh_ else X, O, O, O),
             ("심부전  I50", X, O, O, O),
             ("심근병증  I42~45", X, O, O, O)]
    table(s, Inches(5.9), Inches(1.52), ["질병 (코드)", "허혈성진단비", "심장(특정)", "순환계", "산정특례"], heart, [2.1, 1.05, 0.85, 0.75, 0.75], rh=0.25, fs=8)
    box(s, Inches(5.9), Inches(3.72), Inches(5.3), Inches(0.7),
        "· 빈맥(I47·48) ≠ 부정맥(I49) — 전용 행 구분\n· 순환계 = 2대 + 동맥류·정맥류 등 순환기 전반 [확인]\n● 보장   ○ 미보장", 7.5, False, GREY)
    footer(s, name, 5, TOT)

    # 7 주요치료비 변천사
    s = prs.slides.add_slide(BL); header(s, name, 6, "주요치료비 변천사")
    box(s, Inches(0.42), Inches(1.18), Inches(10), Inches(0.3),
        "3대 주요치료비 담보의 변천사   ①비례형(구간) → ②정액형 → ③비급여 → 생활비", 10.5, True, NAVY)
    for i, (t, d) in enumerate([("1  비례형(구간형)", "치료비 지출액이 구간 하한 도달 시 그 구간 정액 지급.\n2024.11 단종 — 보유자 해지 금지."),
                                ("2  정액형", "치료 사실만으로 약정금 정액 지급.\n면책 90일 · 가입금액 100만원부터."),
                                ("3  비급여(하이클래스)", "급여 전액본인부담금 + 비급여 치료비 정액.\n2천만×10년=2억 · 25.08~ 생활비형 추가."),
                                ("4  왜 필요한가", "5세대 실손: 비중증 비급여 자기부담 50%·한도 1,000만.\n표적·면역 고가 약값엔 부족 → 주요치료비가 메움.")]):
        x = Inches(0.42 + (i % 2) * 5.62); y = Inches(1.56 + (i // 2) * 1.1)
        rect(s, x, y, Inches(5.3), Inches(0.95), LIGHT, LINE)
        box(s, x + Inches(0.1), y + Inches(0.04), Inches(5.1), Inches(0.26), t, 9.5, True, GOLD)
        box(s, x + Inches(0.1), y + Inches(0.3), Inches(5.1), Inches(0.6), d, 8, False, BLACK)
    box(s, Inches(0.42), Inches(3.9), Inches(10), Inches(0.28), "수술 건당 총진료비(평균) · 심평원·건보공단", 10, True, NAVY)
    for i, (k, v) in enumerate([("심장수술(개흉)", "2,832만원"), ("관상동맥우회술", "2,690만원"), ("뇌동맥류 코일색전술", "1,100~1,600만원"),
                                ("심장 스텐트", "1,000~1,400만원"), ("뇌기저부수술", "1,475만원")]):
        x = Inches(0.42 + i * 2.14)
        rect(s, x, Inches(4.22), Inches(1.98), Inches(0.7), WHITE, LINE)
        box(s, x, Inches(4.28), Inches(1.98), Inches(0.26), k, 7.5, False, GREY, PP_ALIGN.CENTER)
        box(s, x, Inches(4.54), Inches(1.98), Inches(0.3), v, 10, True, NAVY, PP_ALIGN.CENTER)
    box(s, Inches(0.42), Inches(5.1), Inches(10.5), Inches(0.5),
        "재발률 — 뇌경색 1년 10%·5년 20~30% / 급성심근경색 1년 30%·3년 50% / 스텐트 5년 재협착 15%", 8.5, False, GREY)
    footer(s, name, 6, TOT)

    # 8 상담 워크시트
    s = prs.slides.add_slide(BL); header(s, name, 7, "상담 워크시트")
    box(s, Inches(0.42), Inches(1.2), Inches(5), Inches(0.28), "암 주요치료비", 11, True, GOLD)
    box(s, Inches(5.9), Inches(1.2), Inches(5), Inches(0.28), "뇌·심장 주요치료비", 11, True, GOLD)
    p5 = {i['t']: i['v'] for i in rep.get('p5_own', [])}
    ch = {c['name']: c['value'] for c in rep.get('chiryo', [])}
    ci_items = {i['t']: i['v'] for i in rep.get('ci', {}).get('items', [])}
    # ★v33b 일반 진단비 + CI 선지급분 병기 (지점장 2026.07.08): 암 5,000만 + ci암 3,200만
    _cov = {c['name']: {i['t']: i['v'] for i in c.get('items', [])} for c in rep.get('coverage', [])}
    _gen_cancer = _cov.get('암', {}).get('일반암')
    def _join(*parts):
        p = [x for x in parts if x and x != '미가입']
        return " + ".join(p) if p else "미가입"
    _cancer = _join(_gen_cancer, (f"ci암 {ci_items['ci암진단비']}" if ci_items.get('ci암진단비') else None))
    # 지침: 뇌·심 진단비 카드 = '뇌혈관 OOO만 · 뇌졸증 OOO만' + CI 선지급분 병기
    bs = " · ".join(f"{k} {v}" for k, v in (("뇌혈관", _cov.get('뇌혈관', {}).get('뇌혈관진단비')),
                                            ("뇌졸증", _cov.get('뇌혈관', {}).get('뇌졸증진단비')))
                    if v and v != '미가입') or "미가입"
    _ci_bs = ci_items.get('ci뇌졸증') or ci_items.get('ci급성심근경색')
    if _ci_bs:
        bs = (bs + f" + ci {_ci_bs}") if bs != "미가입" else f"ci {_ci_bs}"
    cards = [(0.42, "암 진단비", "걸렸을 때 일시금 (기본) + CI 선지급", _cancer),
             (0.42, "암 주요치료비", "수술·방사선·약물 정액 (100만~)", ch.get('암주요치료비', '미가입')),
             (0.42, "하이클래스 (비급여)", "표적·면역·중입자 전액본인 커버", ch.get('비급여주요치료비', '미가입')),
             (5.9, "뇌·심 진단비", "뇌혈관·허혈성 일시금", bs),
             (5.9, "2대 주요치료비", "수술·혈전용해·중환자실 (100만~)", ch.get('순환계주요치료비', '미가입')),
             (5.9, "산정특례 (뇌 · 심장)", "진단만으로 지급 · 개별 담보",
              f"{ch.get('산정특례(뇌혈관)','미가입')} / {ch.get('산정특례(심장)','미가입')}")]
    yl = yr = Inches(1.56)
    for x0, t, sub, val in cards:
        y = yl if x0 == 0.42 else yr; x = Inches(x0)
        ok = '미가입' not in str(val)
        rect(s, x, y, Inches(5.3), Inches(0.95), WHITE, GOLD if ok else LINE)
        box(s, x + Inches(0.12), y + Inches(0.06), Inches(3.1), Inches(0.28), t, 9.5, True, NAVY)
        box(s, x + Inches(0.12), y + Inches(0.34), Inches(3.1), Inches(0.5), sub, 7.5, False, GREY)
        vs = 10.5 if len(str(val)) <= 12 else (8.5 if len(str(val)) <= 26 else 7.5)
        box(s, x + Inches(3.25), y, Inches(1.95), Inches(0.95), str(val), vs, True, NAVY if ok else RED, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
        if x0 == 0.42: yl = y + Inches(1.08)
        else:          yr = y + Inches(1.08)
    rect(s, Inches(0.42), Inches(4.9), W - Inches(0.84), Inches(0.7), LIGHT, GOLD)
    box(s, Inches(0.6), Inches(5.0), Inches(10), Inches(0.5), "부족한 담보 → 보완 추천 :", 10.5, True, NAVY)
    footer(s, name, 7, TOT)

    # 9·10 심혈관 담보 분류
    SETS = [
        [("한화손해보험 4가지", ["심혈관Ⅰ = 협심증 I20 · 기타급성허혈 I24 · 만성허혈 I25 · 빈맥 I47 · 심방세동 I48 · 기타부정맥 I49 · 심부전 I50",
                            "심혈관Ⅱ = 급성심근경색 I21~23 · 인공소생 심장정지 I46.0",
                            "심근병증 = I42 (확장성·비후성·제한성 포함)",
                            "심혈관특정질환 = Ⅰ에서 기타부정맥(I49)만 제외"]),
         ("DB손해보험 4가지", ["특정Ⅰ = 협심증·허혈 I20·24·25 + 주요심장염증 I30~41",
                          "특정Ⅱ = 급성심근경색 I21~23 · I46.0",
                          "특정Ⅲ = 판막질환 · 빈맥 I47 · 심방세동 I48 · 심부전 I50",
                          "순환계 3대 = I46.0 · 부정맥 I47~49 · 심부전 I50"]),
         ("KB손해보험 5가지", ["특정Ⅰ = 협심증·허혈 I20·24·25 · 빈맥 I47·48 · 심부전 I50 (염증X·부정맥X)",
                          "특정Ⅱ = 급성심근경색 I21~23 · I46.0",
                          "심근병증 = I42 / 심장판막질환 = 판막·I33·I38",
                          "기타심장부정맥 = I49 (I47·48은 대상 아님)"]),
         ("현대해상 6가지", ["허혈성 심장질환 = 협심증 I20 · I24 · I25",
                        "특정허혈 심장질환 = 급성심근경색 I21~23",
                        "특정Ⅰ = 빈맥 I47 · 심방세동 I48 · 심부전 I50",
                        "특정Ⅱ = 급성심근경색 I21~23 · I46.0",
                        "주요 심장염증 = I30~41 / 특정2대 = 전도장애 I44·45 + I49",
                        "※ 현대 심근병증 담보 없음"])],
        [("NH농협손해보험 5가지", ["심혈관특정Ⅰ = 협심증·허혈 · 빈맥 · 기타부정맥 I49 · 심부전",
                            "특정Ⅰ(기타부정맥 제외) = Ⅰ에서 I49만 뺀 묶음",
                            "심근병증 = I42·I43 / 주요 심장염증질환 = I30~41",
                            "기타 심장부정맥 = I49"]),
         ("흥국화재 실제 약관", ["특정심혈관질환(기타부정맥제외) = 협심증·허혈·빈맥·심부전 (급성심근 아님)",
                          "특정심혈관질환(기타부정맥) = I49",
                          "심근병증(허혈성제외) = I42·I43 / 주요심장염증질환 = I30~41",
                          "허혈성심질환진단비 = I20·24·25 단독"]),
         ("롯데손해보험 5가지", ["특정심장질환Ⅰ = 급성심근경색 I21~23 · I46.0",
                          "특정심장질환Ⅱ = 협심증·허혈 + 주요심장염증 I30~41",
                          "특정 15대 = 판막 · 심근병증 · 빈맥 · 심부전",
                          "기타 심장부정맥 = I49 / 방실차단·전도장애 = I44·I45"]),
         ("삼성화재 · 메리츠화재", ["허혈성심장질환 6가지 = 급성기 I21 · 후속 I22 · 합병증 I23 · 협심증 I20 · 기타급성 I24 · 만성 I25",
                            "[메리츠] 심장질환진단Ⅰ → 허혈성 / Ⅱ → 급성심근경색 (별도 상품 병존)"])]]
    for si, grp in enumerate(SETS):
        s = prs.slides.add_slide(BL); header(s, name, 8 + si, "심혈관 담보 분류")
        box(s, Inches(0.42), Inches(1.15), W - Inches(0.84), Inches(0.4),
            '★ "특정Ⅰ·Ⅱ"는 회사마다 뜻이 다르다 — 라벨 말고 질병코드로 확인. 빈맥(I47·48)과 부정맥(I49)은 별개.', 8.5, True, RED)
        for i, (comp, lines) in enumerate(grp):
            x = Inches(0.42 + (i % 2) * 5.62); y = Inches(1.62 + (i // 2) * 2.05)
            rect(s, x, y, Inches(5.3), Inches(1.9), WHITE, LINE)
            rect(s, x, y, Inches(5.3), Inches(0.32), NAVY)
            box(s, x, y, Inches(5.3), Inches(0.32), comp, 9.5, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            box(s, x + Inches(0.1), y + Inches(0.36), Inches(5.1), Inches(1.5), "\n".join("· " + l for l in lines), 7.5, False, BLACK)
        footer(s, name, 8 + si, TOT)

    prs.save(out)
    return True


if __name__ == '__main__':
    import sys
    from coverage_benchmark import map_excel_to_report
    if len(sys.argv) < 2:
        print('사용법: python report_pptx.py <엑셀> [고객명]'); sys.exit(1)
    _cn = sys.argv[2] if len(sys.argv) > 2 else '고객'
    _rep = map_excel_to_report(sys.argv[1], settings={'client': _cn, 'branch': '온빛센터 바름지점',
                                                      'manager': '최은혜', 'title': '지점장', 'phone': ''})
    build_report_pptx(_rep, f'보장진단서_{_cn}.pptx')
    print('진단서 PPT 생성 완료(도형·텍스트)')
