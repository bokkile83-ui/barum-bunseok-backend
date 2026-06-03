"""fill_pptx.py v3 — 보장분석 JSON → 보장분석지 PPT 폼 채움.
깨짐 방지: 채운 박스는 폰트 강제(라벨 8pt/버블 길이별 자동축소)+가운데/중앙.
미가입=회색, 가입=진한색, 갱신담보=파랑. 메모 슬라이드는 분류색.
"""
import re
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from datetime import date

GRAY=RGBColor(0xAA,0xAA,0xAA); DARK=RGBColor(0x22,0x22,0x22)
BLUE=RGBColor(0x00,0x00,0xFF); NAVY=RGBColor(0x1F,0x4E,0x79)
ORANGE=RGBColor(0xC0,0x60,0x00); REDC=RGBColor(0xC0,0x00,0x00)

def _n(s): return re.sub(r"[\s()（）·.\-_/]", "", str(s or ""))

def won(m):
    try: m = int(round(float(m)))
    except Exception: return str(m)
    if m <= 0: return "미가입"
    if m >= 10000 and m % 10000 == 0: return f"{m//10000}억"
    if m >= 10000:
        e, t = m // 10000, (m % 10000) // 1000
        return f"{e}억{t}천만" if t else f"{e}억"
    if m >= 1000 and m % 1000 == 0: return f"{m//1000}천만"
    return f"{m}만"

def _sum(cov, subs):
    tot=0; hit=False; renew=False
    for rw in cov:
        nm=_n(rw.get("name"))
        if any(_n(x) in nm for x in subs):
            for a in rw.get("amounts", []):
                try: tot+=int(round(float(a.get("value",0)))); hit=True
                except Exception: pass
                if a.get("renew"): renew=True
    return (tot if hit else None), renew

# (좌표 cm) → 버블/라벨 매핑 — v2와 동일
BUBBLES = {
    (3.6,4.8):["뇌혈관진단"], (3.7,6.3):["뇌졸"], (3.7,7.5):["뇌출혈"],
    (13.5,7.5):["급성심근","허혈성"],
}
LABELBOX = {
    (14.4,15.2):[(["질병 3%","질병3%"],["질병후유3"]),(["질병 80%"],["질병후유80"]),
                 (["상해3%","상해 3%"],["상해후유3"]),(["상해 80%"],["상해후유80"]),(["교통후유"],["교통후유"])],
    (14.4,9.5):[(["대인"],["벌금대인","대인"]),(["대물"],["대물"]),(["합의금"],["처리지원","합의"]),
                (["6주"],["6주"]),(["변호사"],["변호사"]),(["자부상"],["부상위로","자부상"])],
    (14.4,19.5):[(["치매"],["치매","장기요양"]),(["재가"],["재가"]),(["노치원"],["노인전문","노치원"]),(["시설"],["시설급여"])],
    (14.4,24.1):[(["크라운"],["크라운"]),(["임플란트"],["임플란트"]),(["기타"],["치아기타"]),(["화재"],["화재"])],
    (9.9,4.3):[(["협심증","심부전"],["협심증","심부전"]),(["부정맥"],["부정맥"])],
    (1.9,2.3):[(["산정특례"],["산정특례뇌","산정특례(뇌"]),(["혈전용해"],["혈전용해치료비(뇌","혈전용해뇌"]),(["2대주요"],["2대주요치료비(뇌","2대주요뇌"])],
    (11.0,2.2):[(["산정특례"],["산정특례심","산정특례(심"]),(["혈전용해"],["혈전용해치료비(심","혈전용해심"]),(["2대주요"],["2대주요치료비(심","2대주요심"])],
    (0.6,13.9):[(["암주요"],["암주요치료비"])],
    (9.8,13.2):[(["일상","배상"],["일상생활배상"])],
    (9.9,15.6):[(["입원"],["실손입원"]),(["통원"],["실손통원"]),(["MRI"],["MRI"]),(["도수"],["도수"]),(["비급여주사"],["비급여주사"])],
    (9.9,9.4):[(["5대골절"],["5대골절진단비"]),(["골절"],["골절진단비"]),(["중대화상"],["중대화상진단비"]),
               (["화상"],["화상진단비"]),(["반깁스"],["반깁스"]),(["깁스"],["깁스치료"]),(["응급실"],["응급실"])],
    (9.8,19.5):[(["___세","- _","－세"],["질병사망(__"]),(["80세"],["질병사망(80"]),(["종신"],["일반사망(종신","일반사망"])],
    (9.8,21.8):[(["상해사망"],["상해사망"])],
    (0.5,9.3):[(["암진단비"],["암진단비"]),(["유사암"],["유사암"]),(["통합암"],["3대암진단비"]),(["전이암"],["전이암진단비"]),
               (["항암치료"],["항암치료비"]),(["표적"],["표적항암"]),(["양성자"],["양성자"]),(["세기"],["세기조절"]),
               (["다빈치"],["다빈치"]),(["하이클래스"],["하이클래스"]),(["암통원"],["암통원"])],
    (0.3,15.7):[(["질병수술"],["질병수술비"]),(["종합병원수술"],["종합병원수술비(질"]),(["119대"],["119대"]),
                (["뇌혈관수술"],["뇌혈관수술"]),(["허혈성수술","심장수술"],["심장수술"])],
    (4.8,15.8):[(["상해수술"],["상해수술비"]),(["종합병원수술"],["종합병원수술비(상"]),(["5대골절수술"],["5대골절수술"]),
                (["골절수술"],["골절수술"]),(["화상수술"],["화상수술"]),(["중대상해수술"],["중대상해수술"]),(["창상봉합"],["창상봉합"])],
    (0.3,21.9):[(["질병일당"],["질병일당"]),(["상해일당"],["상해일당"]),(["중환자실"],["중환자실"]),
                (["상급병원일당","1인실 상급"],["1인실상급"]),(["종합병원일당","1인실 종합"],["1인실종합"]),
                (["암일당"],["암일당입원","암일당"]),(["간병인"],["간병인사용"]),(["요양병원"],["요양병원"]),(["간호통합"],["간호통합"])],
}

def _match(sh, table, tol=0.6):
    if sh.left is None or sh.top is None: return None
    L,T=Emu(sh.left).cm, Emu(sh.top).cm
    for (x,y),v in table.items():
        if abs(L-x)<tol and abs(T-y)<tol: return v
    return None

def _bubble(sh, text, renew):
    tf=sh.text_frame; tf.word_wrap=False
    try: tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    except Exception: pass
    tf.text=text
    n=len(text)
    fs=11 if n<=2 else 9 if n<=4 else 8 if n<=6 else 7
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size=Pt(fs); r.font.bold=True
        r.font.color.rgb=BLUE if renew else DARK

def _labelfill(sh, mapping, cov):
    tf=sh.text_frame; tf.word_wrap=True
    acc={}
    for rw in cov:
        nm=_n(rw.get("name")); best=-1; blen=0
        for i,(keys,subs) in enumerate(mapping):
            for x in subs:
                xn=_n(x)
                if xn and xn in nm and len(xn)>blen: best,blen=i,len(xn)
        if best>=0:
            a=acc.setdefault(best,[0,False])
            for am in rw.get("amounts",[]):
                try: a[0]+=int(round(float(am.get("value",0))))
                except Exception: pass
                if am.get("renew"): a[1]=True
    lines=[]
    for ln in tf.text.split("\n"):
        if ":" in ln:
            label=ln.split(":")[0]; val=""; renew=False
            for i,(keys,subs) in enumerate(mapping):
                if any(k in label for k in keys):
                    s=acc.get(i)
                    if s and s[0]>0: val=won(s[0]); renew=s[1]
                    else: val="미가입"
                    break
            lines.append((label.rstrip()+" : "+val, val=="미가입", renew))
        else:
            lines.append((ln, False, False))
    tf.text="\n".join(l[0] for l in lines)
    for para,(_,nomiss,renew) in zip(tf.paragraphs, lines):
        for r in para.runs:
            r.font.size=Pt(8)
            r.font.color.rgb=GRAY if nomiss else (BLUE if renew else DARK)

def _walk(shapes, cov):
    for sh in shapes:
        if sh.shape_type==6: _walk(sh.shapes, cov); continue
        if not sh.has_text_frame: continue
        b=_match(sh,BUBBLES)
        if b is not None:
            s,ren=_sum(cov,b); _bubble(sh, won(s) if s else "미가입", ren); continue
        m=_match(sh,LABELBOX)
        if m is not None: _labelfill(sh,m,cov)

def _flatten(shapes):
    for sh in shapes:
        if sh.shape_type==6: yield from _flatten(sh.shapes)
        else: yield sh

def _memo_row(m):
    if isinstance(m, dict):
        return str(m.get("tag","메모")), str(m.get("item","")), str(m.get("note",""))
    t=str(m); tag="메모"
    if any(k in t for k in ("약관확인","확인요망","확인 요망","확인필요","확인 필요")): tag="확인"
    elif "부족" in t: tag="부족"
    elif "미가입" in t: tag="미가입"
    elif "비갱신" in t: tag="비갱신"
    elif "갱신" in t: tag="갱신"
    return tag, "", t

def _tagcolor(tag):
    t=str(tag)
    if "확인" in t: return ("⚠확인", ORANGE)
    if "부족" in t: return ("부족", REDC)
    if "미가입" in t: return ("미가입", GRAY)
    if "비갱신" in t: return ("비갱신", REDC)
    if "갱신" in t: return ("갱신", BLUE)
    return ("메모", DARK)

def _fill_memo_slide(slide, client, memo):
    sh=slide.shapes[0]; tf=sh.text_frame; tf.word_wrap=True
    tf.clear()
    fs=10 if len(memo)<=10 else 8 if len(memo)<=20 else 7
    p0=tf.paragraphs[0]; r0=p0.add_run()
    r0.text=f"셀프팩폭 · 증권확인 ({client})"
    r0.font.size=Pt(fs+3); r0.font.bold=True; r0.font.color.rgb=DARK
    for m in (memo or ["특이사항 없음"]):
        tag,item,note=_memo_row(m); lab,color=_tagcolor(tag)
        p=tf.add_paragraph(); r=p.add_run()
        r.text=f"[{lab}] " + (f"{item} — {note}" if item else note)
        r.font.size=Pt(fs); r.font.color.rgb=color

def fill_pptx(data, xlsx_path, form_path, out_path):
    cov=data.get("coverage",[]); client=data.get("client","고객")
    p=Presentation(form_path)
    _walk(p.slides[0].shapes, cov)
    for sh in p.slides[0].shapes:
        if sh.has_text_frame and "기준" in sh.text_frame.text:
            for r in sh.text_frame.paragraphs[0].runs: pass
            sh.text_frame.paragraphs[0].runs[0].text=f"{date.today():%Y년 %m월 %d일} 기준" if sh.text_frame.paragraphs[0].runs else None
    for sh in _flatten(p.slides[0].shapes):
        if sh.has_text_frame and "보장" in sh.text_frame.text and "님" in sh.text_frame.text:
            if sh.text_frame.paragraphs[0].runs:
                sh.text_frame.paragraphs[0].runs[0].text=f"{client}님의 보장"
    if len(p.slides)>1:
        _fill_memo_slide(p.slides[1], client, data.get("memo",[]))
    p.save(out_path); return out_path
