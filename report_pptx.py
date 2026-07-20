# -*- coding: utf-8 -*-
"""BARUM 보장진단서 PPT v39 — 뼈대 그대로 · 워크시트 흰칸 전부 편집(빈칸 상자 흰색 채우기로 클릭 가능)."""
import os, re, subprocess, tempfile
import xml.etree.ElementTree as ET
from collections import Counter

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NS = {'x': 'http://www.w3.org/1999/xhtml'}
EMU_PER_PT = 12700
DPI = 600   # ★v123 지점장 확정 2026.07.21: 900 → 600dpi (A4 4959x7017).
            #   8K(900)는 11장 래스터에 약 120초가 걸려 분석 1건이 1분을 넘겼다.
            #   600은 약 67초이고 육안·인쇄 품질 차이는 사실상 없다(상업인쇄 표준 300).
            #   구 v89 '150dpi로 낮춰 속도 확보'는 폐기 — 인쇄 시 글자가 뭉개졌다.
            #   래스터 시간은 늘지만 인쇄 품질이 우선. 용량은 아래 적응형 저장으로 방어.
FONT = "맑은 고딕"


def _setfont(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)


def _valueset(rep):
    V = set()

    def add(x):
        if x is None:
            return
        s = str(x).strip().replace(' ', '')
        if s and s not in ('-', '0'):
            V.add(s)

    add(f"{rep.get('n_contract', 0)}건")
    add(f"{rep.get('premium', 0):,}원")
    add(f"{rep.get('renew', 0)}/{rep.get('nonrenew', 0)}")
    add(f"{rep.get('gap_count', 0)}영역")
    for c in rep.get('coverage', []):
        for i in c.get('items', []):
            add(i.get('v'))
    ci = rep.get('ci', {}) or {}
    for i in ci.get('items', []):
        add(i.get('v'))
    add(ci.get('samang')); add(ci.get('residual'))
    if ci.get('rate'):
        add(f"{ci['rate']}%"); add(f"{100 - ci['rate']}%")
    for c in rep.get('chiryo', []):
        add(c.get('value'))
    for d in rep.get('donuts', []):
        add(f"{d.get('pct')}%")
    for d in rep.get('donut_detail', []):
        add(d.get('have')); add(d.get('rec')); add(f"{d.get('pct')}%")
    for c in rep.get('renew_list', []) + rep.get('nonrenew_list', []):
        add(c.get('v'))
    for i in rep.get('p5_own', []):
        add(i.get('v'))
    add(rep.get('client'))
    add(rep.get('band_label'))
    return V


_WS = set()          # 워크시트 흰칸(wbox) 안 문자열


def _patch_worksheet():
    """워크시트 카드의 흰칸 내용을 수집하고, 빈칸엔 '.' 을 넣는다. report_weasy 원본 무수정."""
    import report_weasy as _rw
    if getattr(_rw, '_barum_ws', False):
        return
    _E = '<span class="wbox"></span>'
    _D = '<span class="wbox">.</span>'
    _RX = re.compile(r'<span class="wbox">(.*?)</span>')
    for _fn in ('_wcard', '_wcard_sj', '_wcard_fix'):
        _o = getattr(_rw, _fn, None)
        if _o is None:
            continue

        def _mk(o):
            def _w(*a, **k):
                h = o(*a, **k).replace(_E, _D)
                for m in _RX.finditer(h):
                    t = m.group(1).strip().replace(' ', '')
                    if t:
                        _WS.add(t)
                return h
            return _w
        setattr(_rw, _fn, _mk(_o))

    # ★_wcard_fix_list(.mb 흰칸) 개방: 빈 .mb → '.' 주입, 채워진 값은 _WS 수집
    _EM = '<span class="mb"></span>'
    _DM = '<span class="mb">.</span>'
    _RXM = re.compile(r'<span class="mb">(.*?)</span>')
    def _mkl(o):
        def _w(*a, **k):
            h = o(*a, **k).replace(_EM, _DM)
            for m in _RXM.finditer(h):
                t = m.group(1).strip().replace(' ', '')
                if t and t != '.':
                    _WS.add(t)
            return h
        return _w
    # ★v41 _wcard_fix_group(8p 진단비·수술 그룹칸)도 편집칸으로 개방
    for _fn2 in ('_wcard_fix_list', '_wcard_fix_group'):
        _ol = getattr(_rw, _fn2, None)
        if _ol is not None:
            setattr(_rw, _fn2, _mkl(_ol))

    _rw._barum_ws = True


def _pageset(rep):
    P = set()
    for b in rep.get('premium_bars', []):
        nm = str(b.get('nm', '')).strip().replace(' ', '')
        if nm:
            P.add(nm)
        try:
            P.add(f"{int(b.get('amt', 0)):,}")
        except Exception:
            pass
    # ★v69 하드코딩 인덱스 폐기(2026.07.18): 페이지가 늘거나 순서가 바뀌면
    #   0-based 고정 인덱스(2·5·7·8·9·10)가 통째로 밀려 편집칸이 엉뚱한 페이지에 붙었다.
    #   → _detect_pages()가 PDF 본문에서 표식을 찾아 실제 물리 인덱스를 동적으로 잡는다.
    D = {'.'}
    return {'__bars__': P, '__ws__': set(_WS) | D}


# 편집 대상 페이지를 본문 표식으로 식별(순서·페이지수 변동에 자동 대응)
_PAGE_MARKS = {
    '__bars__': ('AI 진단 요약', '보유 강점', '갱신 / 비갱신 구조'),
    '__ws__':   ('실손보험 세대 구분', '상담 워크시트', '바뀌지 않는 담보',
                 '운전자 · 간병', '재가보험'),
}


def _detect_pages(root, PGSPEC):
    """pdftotext -bbox-layout XML에서 페이지별 텍스트를 읽어 편집대상 물리 인덱스를 찾는다."""
    res = {}
    for pi, pg in enumerate(root.findall('.//x:page', NS)):
        txt = ''.join((w.text or '') for w in pg.findall('.//x:word', NS))
        flat = txt.replace(' ', '')
        for key, marks in _PAGE_MARKS.items():
            if key not in PGSPEC:
                continue
            if any(m.replace(' ', '') in flat for m in marks):
                res.setdefault(pi, set())
                res[pi] |= PGSPEC[key]
    return res


def _value_boxes(xml_path, V, PG=None):
    root = ET.parse(xml_path).getroot()
    PG = _detect_pages(root, PG or {})
    out = []
    for pi, pg in enumerate(root.findall('.//x:page', NS)):
        VV = V | PG.get(pi, set()) | {'.'}
        pw, ph = float(pg.get('width')), float(pg.get('height'))
        found = []
        for ln in pg.findall('.//x:line', NS):
            ws = [(w.text or '', float(w.get('xMin')), float(w.get('yMin')),
                   float(w.get('xMax')), float(w.get('yMax')))
                  for w in ln.findall('x:word', NS)]
            n = len(ws); i = 0
            while i < n:
                for L in (8, 7, 6, 5, 4, 3, 2, 1):
                    if i + L > n:
                        continue
                    grp = ws[i:i + L]
                    if ''.join(g[0] for g in grp) in VV:
                        txt = ' '.join(g[0] for g in grp)
                        found.append((txt, grp[0][1], min(g[2] for g in grp),
                                      grp[-1][3], max(g[4] for g in grp)))
                        i += L
                        break
                else:
                    i += 1
        out.append((pw, ph, found))
    return out


def _fg_of(img, bx):
    x0, y0, x1, y1 = bx
    W, H = img.size
    x0, y0 = max(0, x0), max(0, y0); x1, y1 = min(W, x1), min(H, y1)
    if x1 <= x0 or y1 <= y0:
        return (0, 0, 0)
    ring = []
    for xx in range(x0, x1, max(1, (x1 - x0) // 10)):
        for yy in (max(0, y0 - 2), min(H - 1, y1 + 1)):
            ring.append(img.getpixel((xx, yy))[:3])
    bg = Counter(ring).most_common(1)[0][0] if ring else (255, 255, 255)
    px = list(img.crop((x0, y0, x1, y1)).convert('RGB').getdata())
    return max(px, key=lambda c: (c[0]-bg[0])**2 + (c[1]-bg[1])**2 + (c[2]-bg[2])**2)


def _ink(img, bx):
    x0, y0, x1, y1 = bx
    W, H = img.size
    x0, y0 = max(0, x0), max(0, y0); x1, y1 = min(W, x1), min(H, y1)
    if x1 <= x0 or y1 <= y0:
        return 0.0
    px = list(img.crop((x0, y0, x1, y1)).convert('L').getdata())
    if not px:
        return 0.0
    px.sort(); bg = px[len(px) // 2]
    return sum(1 for v in px if abs(v - bg) > 60) / len(px)


def _erase(img, bx):
    x0, y0, x1, y1 = bx
    W, H = img.size
    x0, y0 = max(0, x0), max(0, y0); x1, y1 = min(W, x1), min(H, y1)
    if x1 <= x0 or y1 <= y0:
        return
    lx, rx = max(0, x0 - 2), min(W - 1, x1 + 1)
    n = x1 - x0
    for yy in range(y0, y1):
        lc = img.getpixel((lx, yy))[:3]
        rc = img.getpixel((rx, yy))[:3]
        if lc == rc:
            for xx in range(x0, x1):
                img.putpixel((xx, yy), lc)
        else:
            for i, xx in enumerate(range(x0, x1)):
                t = i / max(1, n - 1)
                img.putpixel((xx, yy), tuple(int(lc[k] + (rc[k] - lc[k]) * t) for k in range(3)))


def _txt_w(txt, fs):
    """텍스트 렌더 폭(pt) 추정. 한글=fs*1.0, 숫자/영문/기호=fs*0.55, 공백=fs*0.3"""
    w = 0.0
    for ch in txt:
        o = ord(ch)
        if o > 0x1100 and o < 0xD7A4:      # 한글
            w += fs * 1.0
        elif ch == ' ':
            w += fs * 0.3
        elif ch in ',.·/':
            w += fs * 0.35
        else:                               # 숫자·영문·기타
            w += fs * 0.55
    return w



# ══════════════════════════════════════════════════════════════════════════
# ★★★v108 산출물 분할(지점장 확정 2026.07.20, 영구지침)
#   보험 인포메이션 간지 ~ 문서 끝(참고자료 세트) = <PDF·벡터>  → 인쇄해도 선명
#   표지 ~ 재가보험(고객 데이터 페이지)          = <PPT·편집>  → 값 수정용
#   경계는 하드코딩하지 않는다. 계약 7건 이상이면 월보험료 전용 페이지가 끼어
#   간지 위치가 한 장 밀리기 때문(_pageset 고정인덱스 폐기 원칙과 동일).
#   판정 = 간지 고유 문구 '공통 참고 자료'가 있는 물리 페이지.
#   ★v109 지점장 확정: 경계는 <'보험 인포메이션' 글자가 걸리는 페이지>다.
#   앞부분(월보험료 도표 등)이 유동적으로 늘어나면 페이지 번호가 바뀌므로
#   번호가 아니라 <글자>로 잡는다. 숫자 하드코딩 영구 금지.
_INFO_MARK  = '보험인포메이션'      # 1순위: 지점장 지정 문구
_INFO_MARK2 = '참고자료'            # 동반 확인용(간지에만 함께 존재)
_INFO_MARK3 = '공통참고자료'        # 2순위 폴백

def _find_info_page(pdf_path):
    """간지(보험 인포메이션)의 1-based 물리 페이지 번호. 못 찾으면 None.
       판정 = '보험인포메이션' + '참고자료' 동시 → 없으면 '보험인포메이션' 단독 → 없으면 '공통참고자료'."""
    try:
        n = 0
        info = subprocess.run(['pdfinfo', pdf_path], capture_output=True, text=True, timeout=60).stdout
        for ln in info.split('\n'):
            if ln.startswith('Pages:'):
                n = int(ln.split(':')[1].strip()); break
        both = solo = fallback = None
        for i in range(1, n + 1):
            t = subprocess.run(['pdftotext', '-f', str(i), '-l', str(i), pdf_path, '-'],
                               capture_output=True, text=True, timeout=60).stdout
            k = re.sub(r'\s', '', t or '')
            if both is None and (_INFO_MARK in k) and (_INFO_MARK2 in k): both = i
            if solo is None and (_INFO_MARK in k): solo = i
            if fallback is None and (_INFO_MARK3 in k): fallback = i
            if both is not None: break
        return both or solo or fallback
    except Exception:
        pass
    return None


def _pdf_slice(src, dst, first, last):
    """poppler pdfseparate+pdfunite로 first~last 페이지만 잘라 dst에 저장."""
    d = tempfile.mkdtemp()
    pat = os.path.join(d, 'p%d.pdf')
    subprocess.run(['pdfseparate', '-f', str(first), '-l', str(last), src, pat],
                   check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    parts = [os.path.join(d, f'p{i}.pdf') for i in range(first, last + 1)]
    parts = [x for x in parts if os.path.exists(x)]
    if not parts:
        return False
    subprocess.run(['pdfunite'] + parts + [dst], check=True,
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return os.path.exists(dst)

def build_report_pptx(rep, out, dpi=DPI, pdf_out=None):
    """★v107: pdf_out을 주면 내부에서 이미 렌더한 벡터 PDF(보장설명서)를 그 경로로 복사한다.
       렌더 추가 0회 — PPT를 만들 때 어차피 PDF를 한 번 굽기 때문이다.
       PPT는 구조상 'PDF를 사진 찍은 것'이라 확대하면 글자가 뭉갠다. 인쇄용 선명본은 이 PDF다."""
    import report_weasy
    from report_weasy import build_report_pdf
    report_weasy._PPT_MODE = True
    from pdf2image import convert_from_path
    try:
        _WS.clear()
        _patch_worksheet()
    except Exception:
        pass

    tmp = tempfile.mkdtemp()
    pdf = os.path.join(tmp, 'rep.pdf')
    build_report_pdf(rep, pdf)
    # ★v108 분할: 간지 페이지를 찾아 PPT(앞)·PDF(뒤)로 나눈다. 추가 렌더 0회.
    _info = _find_info_page(pdf)
    _npg = 0
    try:
        _pi = subprocess.run(['pdfinfo', pdf], capture_output=True, text=True, timeout=60).stdout
        for _ln in _pi.split('\n'):
            if _ln.startswith('Pages:'):
                _npg = int(_ln.split(':')[1].strip()); break
    except Exception:
        pass
    if pdf_out:
        try:
            import shutil as _sh
            if _info and _npg and _info <= _npg:
                if not _pdf_slice(pdf, pdf_out, _info, _npg):
                    _sh.copyfile(pdf, pdf_out)      # 분할 실패 → 통본 폴백(누락 금지)
            else:
                _sh.copyfile(pdf, pdf_out)          # 간지 미검출 → 통본 폴백
        except Exception:
            try:
                import shutil as _sh2; _sh2.copyfile(pdf, pdf_out)
            except Exception:
                pass

    try:
        xml = os.path.join(tmp, 'bb.xml')
        subprocess.run(['pdftotext', '-bbox-layout', pdf, xml], check=True,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        pages = _value_boxes(xml, _valueset(rep), _pageset(rep))
    except Exception:
        pages = []

    # ★v108: PPT는 간지 <앞>까지만 굽는다(고객 데이터 페이지). 뒤는 PDF가 담당.
    #   래스터 장수가 29→11로 줄어 300dpi로 올려도 오히려 기존보다 빠르다.
    _last = (_info - 1) if (_info and _info > 1) else None
    if not _last:
        _last = 0
        try:
            _pi2 = subprocess.run(['pdfinfo', pdf], capture_output=True, text=True, timeout=60).stdout
            for _ln in _pi2.split('\n'):
                if _ln.startswith('Pages:'):
                    _last = int(_ln.split(':')[1].strip()); break
        except Exception:
            _last = len(pages) or 1
    # ★★v112 메모리 방어(8K 필수조건): 전 페이지를 한꺼번에 변환하면
    #   900dpi에서 장당 약 224MB × 11장 = 2.4GB가 동시에 RAM에 뜬다 → Railway OOM.
    #   그래서 <한 장씩 변환 → 저장 → 즉시 해제>로 바꾼다. 순간 점유 = 1장분(약 224MB).
    _n_pages = _last

    prs = Presentation()
    pw_pt, ph_pt = (pages[0][0], pages[0][1]) if pages else (595.276, 841.890)
    prs.slide_width = Emu(int(pw_pt * EMU_PER_PT))
    prs.slide_height = Emu(int(ph_pt * EMU_PER_PT))
    BL = prs.slide_layouts[6]

    for idx in range(_n_pages):
        im = convert_from_path(pdf, dpi=dpi, first_page=idx + 1, last_page=idx + 1)[0]
        im = im.convert('RGB')
        vals = pages[idx][2] if idx < len(pages) else []
        sx = im.width / pw_pt
        sy = im.height / ph_pt

        meta = []
        for txt, x0, y0, x1, y1 in vals:
            bx = (int(x0 * sx) - 1, int(y0 * sy) - 1, int(x1 * sx) + 2, int(y1 * sy) + 1)
            meta.append((txt, x0, y0, x1, y1, _fg_of(im, bx), _ink(im, bx), bx))
        for *_, bx in meta:
            _erase(im, bx)

        # ★v106 적응형 저장: 문서형 페이지는 PNG가 더 작고 선명하다.
        #   사진·이미지형 페이지(재무 3장 등)는 PNG가 급격히 커지므로 1.5MB 초과 시 고품질 JPEG로 대체.
        ip = os.path.join(tmp, f'p{idx}.png')
        im.save(ip, 'PNG')
        try:
            if os.path.getsize(ip) > 1_500_000:
                ij = os.path.join(tmp, f'p{idx}.jpg')
                im.save(ij, 'JPEG', quality=94, subsampling=0, optimize=True)
                if os.path.getsize(ij) < os.path.getsize(ip):
                    os.unlink(ip); ip = ij
        except Exception:
            pass

        s = prs.slides.add_slide(BL)
        s.shapes.add_picture(ip, 0, 0, prs.slide_width, prs.slide_height)
        try:
            im.close(); del im      # ★v112: 900dpi 1장분(약 224MB) 즉시 해제 — 다음 장 변환 전에 비운다
        except Exception:
            pass

        for txt, x0, y0, x1, y1, fg, ink, _bx in meta:
            h_pt = y1 - y0
            fs = max(5.0, round(h_pt * 0.90, 1))
            pad = 2.0
            if txt == '.':
                x0 = x1 - 28.0
                fs = 8.5
            else:
                # ★폭 기반 자동 축소: 문자열이 상자폭 넘으면 폰트 줄임(오버 방지)
                box_w = (x1 - x0) + pad * 1.0   # 가용 폭(pt), 좌우 여백 약간 남김
                need = _txt_w(txt, fs)
                if need > box_w and need > 0:
                    fs = max(5.0, round(fs * box_w / need, 1))
            sh = s.shapes.add_textbox(Emu(int((x0 - pad) * EMU_PER_PT)),
                                      Emu(int((y0 - pad * 0.6) * EMU_PER_PT)),
                                      Emu(int((x1 - x0 + pad * 2.4) * EMU_PER_PT)),
                                      Emu(int((h_pt + pad * 1.2) * EMU_PER_PT)))
            tf = sh.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if txt == '.' else PP_ALIGN.LEFT
            r = p.add_run(); r.text = txt
            r.font.size = Pt(fs)
            r.font.color.rgb = RGBColor(*fg)
            r.font.bold = (ink > 0.30)
            _setfont(r)
            if txt == '.':
                f = sh.fill
                f.solid()
                f.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                sh.line.fill.background()

    prs.save(out)
    return True
