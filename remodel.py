# ===== BARUM remodel.py v526-grp-20260821 =====
# ★★★★★ 보험 리모델링 비교 (지점장 지시 2026.08.15)
#   지점장 원문: 「새앱을 만들자 / 1버튼 기존 보험 엑셀 / 2버튼 새로 정리된 엑셀
#                → 그럼 두개를 비교한 진단서 / 틀은 저걸로」
#              「수정을 생각하면 엑셀, 이미지를 생각하면 ppt」 → <b>산출물 = PPT</b>
#
# ★설계 원칙
#   1. <b>기존 `analyze()` 파이프라인을 일절 건드리지 않는다</b>(회귀 위험 0).
#   2. 입력은 <b>우리 앱이 만든 보장분석지 xlsx 2개</b>다. PDF 파싱이 없다.
#   3. 담보명은 master.xlsx 101행으로 <b>이미 고정</b>이라 두 파일이 1:1로 대조된다.
#   4. <b>삭제 특약은 물어보지 않는다</b> — 기존에 있고 새 파일에 없으면 그것이 삭제다(차집합).
#
# ★엑셀 읽기 규격 (build_excel 산출물)
#   1행 = 회사\n상품\n[갱신] · 2행 = 월보험료 · 6행~ = 담보(B열 담보명)
#   값은 <b>끝열(합계)</b>에서 읽는다. 단 헤더가 '보유 합계'·'제안 합계'·'합계'인 열은
#   계약이 아니다(v419 조문) — 보험료 합계를 낼 때 제외한다.

import re, io

_SUMHDR = ('보유 합계', '제안 합계', '합계')


def _lump(v):
    """★일시납 금액은 헤더 2행 「11,000,000 (일시납)」 텍스트에 있다.
       _num은 0으로 보지만(월보험료 합계 방어) 자산 페이지에는 실어야 한다."""
    import re as _re
    t = str(v or '')
    if '일시납' not in t:
        return 0
    m = _re.search(r'([\d,]+)', t)
    return int(m.group(1).replace(',', '')) if m else 0


def _num(v):
    """셀값 → 숫자. 슬래시 표기(20/50/100)는 <b>최댓값</b>으로 본다(대표값)."""
    if v is None:
        return 0.0
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip()
    if not s or s.startswith('='):
        return 0.0
    # ★★★★★v424 (실측 2026.08.16): 헤더 2행의 「11,000,000 (일시납)」이 월보험료에 더해져
    #   합계가 386,530 → <b>11,386,530</b>이 됐다. 엑셀 SUM은 텍스트를 무시하는데
    #   이 파서는 숫자를 뽑아 썼다. <b>일시납·완납 표기는 월보험료가 아니다</b> → 0으로 본다.
    if ('일시납' in s) or ('완납' in s):
        return 0.0
    if '/' in s:
        ns = [float(re.sub(r'[^\d.]', '', p) or 0) for p in s.split('/')]
        return max(ns) if ns else 0.0
    s = re.sub(r'[^\d.\-]', '', s)
    try:
        return float(s) if s else 0.0
    except Exception:
        return 0.0


def read_sheet(path_or_bytes):
    """보장분석지 xlsx → {'premium':월보험료합, 'contracts':[회사·상품], 'cov':{담보명:금액}}"""
    import openpyxl
    src = path_or_bytes
    if isinstance(src, (bytes, bytearray)):
        src = io.BytesIO(src)
    wb = openpyxl.load_workbook(src, data_only=True)
    ws = wb['보장분석'] if '보장분석' in wb.sheetnames else wb[wb.sheetnames[0]]

    # 계약 열 = 헤더가 합산 라벨이 아닌 열 (v419: 「계약이냐」 판정은 한 곳에서)
    ct_cols, sum_cols = [], []
    for c in range(3, ws.max_column + 1):
        h = str(ws.cell(1, c).value or '').strip()
        if not h:
            continue
        (sum_cols if h in _SUMHDR else ct_cols).append(c)

    premium = 0.0
    contracts = []
    for c in ct_cols:
        h = str(ws.cell(1, c).value or '')
        parts = [p.strip() for p in h.split('\n') if p.strip()]
        contracts.append({'company': parts[0] if parts else '',
                          'product': parts[1] if len(parts) > 1 else '',
                          'renewal': (parts[2] if len(parts) > 2 else '').strip('[]'),
                          'lump_sum': _lump(ws.cell(2, c).value),
                          # ★★★★★v463 제71조 2항 (지점장 지적 2026.08.17 「리포트에도 종신+연금 나오는지」)
                          #   실측: 종신 판별은 맞았는데 표의 <b>가입날짜·납입기간 칸이 '—'</b>였다.
                          #   read_sheet가 두 키를 안 만들었다. 마스터 헤더 3행=가입년일 · 5행=총납입기간.
                          'contract_date': str(ws.cell(3, c).value or '').strip(),
                          # ★★★★★v530 제71조 3항 (지점장 지시 2026.08.21 「진단서·리포트에
                          #   종신·연금 안 나오는 것도 점검」) — 이 경로에만 <b>만기(4행)가 없었다</b>.
                          #   종신 판별의 정본은 §8.1 <b>만기 9999</b>인데 키가 없어 상품명 글자로만 갈렸다.
                          'expiry_date': str(ws.cell(4, c).value or '').strip(),
                          'pay_term': str(ws.cell(5, c).value or '').strip(),
                          'premium': _num(ws.cell(2, c).value)})
        premium += _num(ws.cell(2, c).value)

    # 값은 끝열(합계)에서. 합산 열이 없으면 계약 열을 더한다.
    endc = sum_cols[-1] if sum_cols else None
    cov = {}
    for r in range(6, ws.max_row + 1):
        nm = ws.cell(r, 2).value
        if not nm:
            continue
        nm = str(nm).strip()
        if endc:
            v = _num(ws.cell(r, endc).value)
            if not v:   # =SUM 캐시가 없으면 계약 열 직접 합산(제2조 등식2)
                v = sum(_num(ws.cell(r, c).value) for c in ct_cols)
        else:
            v = sum(_num(ws.cell(r, c).value) for c in ct_cols)
        cov[nm] = v
    return {'premium': premium, 'contracts': contracts, 'cov': cov}


def split_sheet(path_or_bytes):
    """★★★★★v422h — 지점장 확정 2026.08.15 「엑셀은 1번만 보면된다」

    엑셀 <b>한 개</b> 안에 이미 기존과 최종이 둘 다 들어 있다(v388c 합산 2열).
      ㆍ기존   = 보유 계약 열들 (헤더 `제안 합계` 바로 왼쪽 열은 제안 계약 — v419)
      ㆍ최종   = 보유 + 제안
    두 파일을 받지 않는다. 반환은 `read_sheet`와 <b>같은 모양</b>이라 `compare`가 그대로 돈다.
    """
    import openpyxl
    src = path_or_bytes
    if isinstance(src, (bytes, bytearray)):
        src = io.BytesIO(src)
    wb = openpyxl.load_workbook(src, data_only=True)
    ws = wb['보장분석'] if '보장분석' in wb.sheetnames else wb[wb.sheetnames[0]]

    ct_cols, sum_cols, hdr = [], {}, {}
    for c in range(3, ws.max_column + 1):
        h = str(ws.cell(1, c).value or '').strip()
        if not h:
            continue
        hdr[c] = h
        if h in _SUMHDR:
            sum_cols[h] = c
        else:
            ct_cols.append(c)

    # ★제안 계약 열은 <b>엑셀이 스스로 말한다</b> — `제안 합계` 열의 `=SUM(J6:J6)` 범위가 곧 제안 구간이다.
    #   위치·순서를 추측하지 않는다(제11조 구조 가정 금지). 수식이 없으면 주황 헤더(ED7D31)로 폴백.
    prop_cols = []
    if '제안 합계' in sum_cols:
        pc = sum_cols['제안 합계']
        wf = openpyxl.load_workbook(io.BytesIO(path_or_bytes) if isinstance(path_or_bytes, (bytes, bytearray))
                                    else path_or_bytes, data_only=False)
        wsf = wf['보장분석'] if '보장분석' in wf.sheetnames else wf[wf.sheetnames[0]]
        for r in range(6, wsf.max_row + 1):
            f = str(wsf.cell(r, pc).value or '')
            m = re.match(r'^=SUM\(([A-Z]+)\d+:([A-Z]+)\d+\)$', f.replace(' ', ''))
            if m:
                from openpyxl.utils import column_index_from_string as _ci
                a, b = _ci(m.group(1)), _ci(m.group(2))
                prop_cols = [c for c in ct_cols if a <= c <= b]
                break
        if not prop_cols:
            for c in ct_cols:
                rgb = str(getattr(wsf.cell(1, c).fill.fgColor, 'rgb', '') or '')
                if rgb.upper().endswith('ED7D31'):
                    prop_cols.append(c)
    own_cols = [c for c in ct_cols if c not in prop_cols]

    def _mk(cols):
        cs, prem = [], 0.0
        for c in cols:
            parts = [p.strip() for p in str(ws.cell(1, c).value or '').split('\n') if p.strip()]
            p = _num(ws.cell(2, c).value)
            # ★v424 일시납 금액 · 가입날짜 · 납입기간을 함께 담는다(가입금액이 곧 담보)
            _h2 = str(ws.cell(2, c).value or '')
            _lp = 0
            if '일시납' in _h2:
                _mm = re.search(r'([\d,]+)', _h2)
                if _mm:
                    _lp = int(_mm.group(1).replace(',', ''))
            cs.append({'company': parts[0] if parts else '',
                       'product': parts[1] if len(parts) > 1 else '',
                       'renewal': (parts[2] if len(parts) > 2 else '').strip('[]'),
                       'lump_sum': _lp,
                       'contract_date': ws.cell(3, c).value or '',
                       'expiry_date': ws.cell(4, c).value or '',
                       'pay_term': ws.cell(5, c).value or '',
                       'premium': p})
            prem += p
        return cs, prem

    own_ct, own_prem = _mk(own_cols)
    prop_ct, prop_prem = _mk(prop_cols)

    # ★★★★★v422u — 담보는 <b>행 단위</b>로 잡는다(지점장 「엑셀의 담보는 고정이다」 2026.08.15).
    #   담보명을 dict 키로 쓰면 <b>동명 담보 2행</b>(2대 주요치료비 · 혈전용해치료비)이 1행으로 합쳐진다.
    #   101행 → 99행이 되는 것은 마스터를 재정렬한 것과 같다.
    cov_old, cov_new, rows = {}, {}, []
    _grp = ''
    for r in range(6, ws.max_row + 1):
        # ★구분(A열)은 그룹 첫 행에만 있다 — 이어받는다(지점장 「암·뇌·심장처럼 앞에」 2026.08.15)
        _g = ws.cell(r, 1).value
        if _g and str(_g).strip():
            _grp = str(_g).strip()
        nm = ws.cell(r, 2).value
        if not nm:
            continue
        nm = str(nm).strip()
        o = sum(_num(ws.cell(r, c).value) for c in own_cols)
        n = o + sum(_num(ws.cell(r, c).value) for c in prop_cols)
        rows.append((r, nm, o, n, _grp))
        cov_old[nm] = cov_old.get(nm, 0.0) + o
        cov_new[nm] = cov_new.get(nm, 0.0) + n

    old = {'premium': own_prem, 'contracts': own_ct, 'cov': cov_old, 'rows': rows}
    new = {'premium': own_prem + prop_prem, 'contracts': own_ct + prop_ct, 'cov': cov_new, 'rows': rows}
    return old, new, bool(prop_cols)


def contract_kinds(contracts):
    """계약 목록에서 연금 · 종신 · 저축을 가려낸다.
       담보 금액 체계(가입금액)와 성질이 달라 <b>마스터 행으로 만들지 않는다</b>."""
    out = {'연금': [], '종신': [], '저축': []}
    for c in contracts or []:
        p = str(c.get('product') or '')
        r = str(c.get('renewal') or '')
        if '연금' in p:
            out['연금'].append(c)
        # ★v530 제71조 3항 — §8.1 정본은 <b>만기 9999(종신)</b>다. 상품명에 「종신」이
        #   없는 종신계약(실측 교보 `교보3밸런스보장보험`)이 통째로 빠져 있었다.
        _ex = str(c.get('expiry_date') or '')
        if '종신' in p or '종신' in r or '9999' in _ex:
            out['종신'].append(c)
        if '저축' in p:
            out['저축'].append(c)
    return out


def _rowlist(old, kind):
    """★행 단위 분류 — 담보명 dict는 동명 2행을 합쳐 값을 두 배로 만든다."""
    # ★값은 <b>행 단위</b>로 읽되(합산 금지), 같은 담보가 두 행이면 <b>대표값 한 줄</b>로 접는다.
    #   고객 눈에 같은 이름이 두 번 나오면 그건 오류로 보인다.
    # ★★★★★v427 (실측 2026.08.16): `rows`는 <b>엑셀 1개 경로</b>(split_sheet)에서만 만들어진다.
    #   <b>엑셀 2개 경로</b>(read_sheet ×2)에는 없어 이 함수가 빈 목록을 돌려주고
    #   삭제·감소가 <b>영원히 0</b>이 됐다. rows가 없으면 호출부의 dict 계산을 그대로 쓴다.
    if not old.get('rows'):
        return None
    best = {}
    order = []
    for _r, nm, o, n, _g in old.get('rows', []):
        if o == 0 and n == 0:
            continue
        k = ('add' if o == 0 else 'delete' if n == 0 else
             'up' if n > o else 'down' if n < o else 'same')
        if k != kind:
            continue
        if nm not in best:
            best[nm] = (nm, o, n, n - o); order.append(nm)
        elif abs(n - o) > abs(best[nm][3]) or n > best[nm][2]:
            best[nm] = (nm, o, n, n - o)
    return [best[x] for x in order]



_GRPMAP = None


def _grp_of(nm):
    """담보명 → 구분(A열). 마스터를 한 번만 읽어 캐시한다.
       ★구분을 코드에 적지 않는다 — 마스터가 정본이다(제15조)."""
    global _GRPMAP
    if _GRPMAP is None:
        _GRPMAP = {}
        try:
            import openpyxl, os
            _p = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'master.xlsx')
            _ws = openpyxl.load_workbook(_p)['보장분석']
            _cur = ''
            for _r in range(6, _ws.max_row + 1):
                _a = _ws.cell(_r, 1).value
                if _a and str(_a).strip():
                    _cur = str(_a).strip()
                _b = _ws.cell(_r, 2).value
                if _b:
                    _GRPMAP[str(_b).strip()] = _cur
        except Exception as _e:
            print('[v468 구분] 마스터 읽기 실패', str(_e)[:60])
    return _GRPMAP.get(str(nm).strip(), '기타')


def compare(old, new):
    """기존 ↔ 신규 비교 결과."""
    names = list(old['cov'].keys())
    for k in new['cov']:
        if k not in names:
            names.append(k)

    # ★★★★★v422t — 「엑셀의 담보는 고정이다」(지점장 2026.08.15).
    #   담보 행은 <b>엑셀 순서 그대로</b> 전부 싣는다. 변화가 없다고 빼면 그건 마스터를 재정렬한 것이다.
    #   미가입(0/0)도 엑셀에 행이 있으므로 뺀다 = 행 삭제다. 뺄 수 없다.
    up, down, same, dele, add, allrows = [], [], [], [], [], []
    for nm in names:
        o = old['cov'].get(nm, 0.0)
        n = new['cov'].get(nm, 0.0)
        # ★★★★★v468 제75조 4항 (지점장 2026.08.17 「말은 101가지 담보라면서 엑셀도 1페이지다」)
        #   ★결함: `allrows`가 <b>비어 있었다</b>. 담보표 원천이 `old['rows']`뿐인데
        #   `read_sheet`는 `rows`를 만들지 않는다(키: contracts·cov·premium 뿐 — 실측).
        #   그래서 담보표가 통째로 0행 → 엑셀이 1쪽이었다.
        #   → <b>담보 전수</b>를 여기서 만든다. 미가입(0/0)도 싣는다 —
        #     「없다」는 것도 상담에서 보여줄 정보다(지점장 시안도 전 담보를 싣는다).
        allrows.append((_grp_of(nm), nm, o, n, n - o,
                        '미가입' if (o == 0 and n == 0) else
                        '삭제' if n == 0 else '신규 추가' if o == 0 else
                        '보장 증가' if n > o else '보장 감소' if n < o else '변동 없음'))
        if o == 0 and n == 0:
            continue
        if o > 0 and n == 0:
            dele.append((nm, o, n, n - o))
        elif o == 0 and n > 0:
            add.append((nm, o, n, n - o))
        elif n > o:
            up.append((nm, o, n, n - o))
        elif n < o:
            down.append((nm, o, n, n - o))
        else:
            same.append((nm, o, n, 0.0))

    up.sort(key=lambda x: -x[3])
    add.sort(key=lambda x: -x[2])
    dele.sort(key=lambda x: -x[1])
    down.sort(key=lambda x: x[3])

    # ★★★★★v422e (지점장 지시 2026.08.15 「이건 엑셀로주고 이건 리포트로 ppt로」)
    #   틀 사진의 <b>「제안 보험료 + 삭제 후 보험료」</b> 분해도 <b>입력 없이 자동</b>이다.
    #   새 엑셀의 계약 중 <b>기존 엑셀에도 있던 회사·상품</b> = 유지(삭제 후 잔존),
    #   <b>새로 생긴 것</b> = 제안. 「무엇을 해지했나」를 사람에게 묻지 않는다.
    _okey = {(c['company'], c['product']) for c in old['contracts']}
    keep, prop = [], []
    for c in new['contracts']:
        (keep if (c['company'], c['product']) in _okey else prop).append(c)
    kill = [c for c in old['contracts']
            if (c['company'], c['product']) not in {(x['company'], x['product']) for x in new['contracts']}]

    op, np_ = old['premium'], new['premium']
    return {'keep': keep, 'prop': prop, 'kill': kill,
            'prem_keep': sum(c['premium'] for c in keep),
            'prem_prop': sum(c['premium'] for c in prop),
            'old': old, 'new': new,
            'prem_old': op, 'prem_new': np_,
            'save_m': op - np_, 'save_y': (op - np_) * 12,
            'save_pct': (round((op - np_) / op * 100, 1) if op else 0.0),
            'up': _rowlist(old, 'up') or up, 'down': _rowlist(old, 'down') or down,
            'same': _rowlist(old, 'same') or same, 'delete': _rowlist(old, 'delete') or dele,
            'add': _rowlist(old, 'add') or add,
            'all': [(_g, nm, o, n, n - o,
                     '미가입' if (o == 0 and n == 0) else
                     '삭제' if n == 0 else '신규 추가' if o == 0 else
                     '보장 증가' if n > o else '보장 감소' if n < o else '변동 없음')
                    for _r, nm, o, n, _g in old.get('rows', [])] or allrows}


# ─────────────────────────────── PPT ───────────────────────────────
def _mw(v):
    """만원 표기."""
    try:
        v = int(round(float(v)))
    except Exception:
        return '-'
    return f'{v:,}만원'

def build_report(cmp_, client="고객", base_date="", total=9):
    """★★★★★v423 — 리포트 7쪽을 <b>지점장 시안 HTML</b>로 만들어 (PDF, [PNG]) 로 돌려준다.
       구판 `build_pptx`(PPT 도형)는 폐기했다. 뷰어가 차트를 자기 방식으로 다시 그려
       데이터 레이블 서식도 막대 색도 버렸다 — 내가 고칠 수 없는 자리였다."""
    import tempfile, os
    import report_pages
    from weasyprint import HTML
    from pypdf import PdfWriter
    from pdf2image import convert_from_path

    tmp = tempfile.mkdtemp()
    pdfs, pngs = [], []
    import os as _os
    _base = _os.path.dirname(_os.path.abspath(report_pages.__file__)) + '/'
    # ★★★★★v530 제121조 — 발행 직전 <b>심장 값 대조</b>. 엑셀과 다르면 그 자리에서 멈춘다.
    #   지점장 지시 2026.08.21 「네 산출물의 담보값을 실제로 대조해 다르면 발행을 막는 검사」
    # ★★★★★v573 긴급 (지점장 실측 2026.08.23 「엑셀 2개를 넣는데 심장동결이 계속 떠서
    #   비교가 안 된다 · 아예 최종 엑셀 업로드 자체가 안 된다 · 심장동결은 기존보험 유지로
    #   인해 가능한 일이다 · 케이스 4가지는 별개다」).
    #   [원인] `heart_audit`는 <b>`cmp_['new']['cov']` 하나만</b> 본다. 그런데 리모델링 리포트는
    #     제127조로 <b>보유(검정) + 제안 증가분(레드)</b>을 함께 찍는다 — <b>기존 유지분이 더해져</b>
    #     엑셀 500 / 리포트 1,403 처럼 어긋난다. <b>불일치가 아니라 설계다.</b>
    #   [범위] 제121조는 <b>진단서 4·5쪽</b>(analyze 경로) 조문이다. 리모델링 리포트는 <b>다른 산출물</b>이다.
    #   ⇒ 리모델링에서는 <b>차단하지 않는다.</b> 대신 로그로 남긴다(조용히 틀리지 않는다).
    #     진단서 차단(main.py 제121조 게이트)은 <b>그대로 둔다</b>.
    _hb = report_pages.heart_audit(cmp_, client, base_date)
    if _hb:
        print('[제121조 심장] 리모델링 리포트 — 보유+증가분 표기로 차이 %d건 (차단하지 않음)' % len(_hb))
        for _x in _hb[:6]: print('   ·', _x)
    else:
        print('[제121조 심장동결] 리포트 4·5쪽 = 엑셀 · 불일치 0건')

    for i, html in enumerate(report_pages.build(cmp_, client, base_date, total), 1):
        f = os.path.join(tmp, 'p%d.pdf' % i)
        HTML(string=html, base_url=_base).write_pdf(f)
        # ★★★★★v483 제93조 (지점장 2026.08.19 「출력하면 희미해서 보기가 싫다 · 진하게 해줘」)
        #   200dpi도 인쇄에서 흐렸다 → <b>300dpi</b>. 본문 회색 글자도 함께 진하게(report_pages).
        ims = convert_from_path(f, dpi=300)
        if len(ims) != 1:                      # ★한 쪽이 넘치면 조용히 넘어가지 않는다
            print('[REPORT_OVERFLOW] %d쪽이 %d장이 됐다' % (i, len(ims)))
        g = os.path.join(tmp, 'p%d.png' % i)
        ims[0].save(g)
        pdfs.append(f); pngs.append(g)

    w = PdfWriter()
    for f in pdfs:
        w.append(f)
    out = os.path.join(tmp, 'report.pdf')
    w.write(out); w.close()
    return open(out, 'rb').read(), pngs


# ★★★★★v546 제129조 3항 (지점장 지시 2026.08.22 「입력값 넣는곳은 다 수정이 되도록 · 혹시나 이름도」)
#   전체 텍스트화(v543)는 파워포인트에서 깨졌다 — 651개 상자가 서로 밀렸다.
#   ⇒ <b>값이 들어가는 칸만</b> 텍스트로 만든다. <b>숫자가 든 상자 + 고객명</b>만 얹고
#     제목·라벨·설명은 <b>이미지 그대로</b> 둔다. 상자 수가 1/10로 줄어 밀림이 없다.
_TEXT_PPT = True         # ★v546 — 값 칸만 텍스트(제129조 3항)


def build_report_pptx(pngs, pdf_bytes=None, client=''):
    """★★★★★v543 제129조 (지점장 지시 2026.08.22 「여전히 리포트ppt는 통으로된 이미지다」)
       구 코드는 PNG를 <b>그냥 붙이기만</b> 했다 — 슬라이드마다 도형 1개(그림)·텍스트 0개.
       ⇒ 리포트 PDF의 <b>모든 글자를 좌표대로 뽑아</b> 배경에서 지우고 <b>텍스트 상자로 얹는다</b>.
         배경 PNG에는 선·표·색만 남고, 글자는 전부 <b>고를 수 있는 텍스트</b>가 된다.
       pdf_bytes가 없거나 추출이 실패하면 <b>종전대로 통이미지</b>로 낸다(산출물 누락 금지)."""
    from pptx import Presentation
    from pptx.util import Cm, Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(21.0), Cm(29.7)
    bl = prs.slide_layouts[6]

    # ★★★★★v544 제129조 2항 (지점장 실측 2026.08.22 「지금 보이는 이미지 실화냐? 다 깨졌는데」)
    #   내 검증은 <b>LibreOffice 렌더</b>였고 지점장은 <b>파워포인트</b>로 연다 — 렌더러가 다르면
    #   글자 폭·줄바꿈이 달라 <b>651개 상자가 서로 밀려 깨진다</b>.
    #   ⇒ <b>기본은 통이미지</b>로 되돌린다(작동 보장). 텍스트화는 `_TEXT_PPT=True`일 때만 켠다.
    #   ★「고쳤다」는 <b>지점장이 쓰는 프로그램에서 확인한 뒤</b>에만 쓴다. LibreOffice는 대역이 아니다.
    boxes = []
    if pdf_bytes and _TEXT_PPT:
        try:
            import subprocess, tempfile, os as _os
            import report_pptx as _rp
            _td = tempfile.mkdtemp()
            _pf = _os.path.join(_td, 'r.pdf'); open(_pf, 'wb').write(pdf_bytes)
            _xf = _os.path.join(_td, 'r.xml')
            subprocess.run(['pdftotext', '-bbox-layout', _pf, _xf], check=True,
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            _raw = _rp._all_boxes(_xf)
            # ★v546 제129조 3항 — <b>값 칸만</b> 고른다: 숫자가 든 상자 + 고객명.
            #   라벨·제목·설명은 이미지 그대로 두어 밀림을 막는다.
            import re as _re9
            _cl = str(client or '').strip()
            boxes = []
            for _pw, _ph, _ws in _raw:
                _keep = [w for w in _ws
                         if _re9.search(r'\d', w[0]) or (_cl and _cl in w[0])]
                boxes.append((_pw, _ph, _keep))
            print('[v546 리포트PPT] 값칸 텍스트 %d쪽 · 상자 %d개(전체 %d개 중)'
                  % (len(boxes), sum(len(b[2]) for b in boxes),
                     sum(len(b[2]) for b in _raw)))
        except Exception as e:
            print('[v543 리포트PPT] 텍스트 추출 실패 → 통이미지 폴백:', e)
            boxes = []

    from PIL import Image
    import report_pptx as _rp2
    EMU_PT = 12700
    for i, g in enumerate(pngs):
        s = prs.slides.add_slide(bl)
        if i < len(boxes) and boxes[i][2]:
            pw, ph, words = boxes[i]
            im = Image.open(g).convert('RGB')
            sx, sy = im.width / pw, im.height / ph
            meta = []
            for txt, x0, y0, x1, y1 in words:
                bx = (int(x0 * sx) - 1, int(y0 * sy) - 1, int(x1 * sx) + 2, int(y1 * sy) + 1)
                meta.append((txt, x0, y0, x1, y1, _rp2._fg_of(im, bx), _rp2._ink(im, bx), bx))
            for *_r, bx in meta:
                _rp2._erase(im, bx)
            _tp = g + '.txt.png'; im.save(_tp, 'PNG'); im.close()
            s.shapes.add_picture(_tp, Cm(0), Cm(0), Cm(21.0), Cm(29.7))
            for txt, x0, y0, x1, y1, fg, ink, _bx in meta:
                h_pt = y1 - y0
                fs = max(4.5, round(h_pt * 0.88, 1))
                _need = _rp2._txt_w(txt, fs); _bw = (x1 - x0) + 2.0
                if _need > _bw and _need > 0:
                    fs = max(4.5, round(fs * _bw / _need, 1))
                sh = s.shapes.add_textbox(Emu(int((x0 - 1.6) * EMU_PT)), Emu(int((y0 - 1.2) * EMU_PT)),
                                          Emu(int((x1 - x0 + 3.6) * EMU_PT)), Emu(int((h_pt + 2.4) * EMU_PT)))
                tf = sh.text_frame; tf.word_wrap = False
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
                r = p.add_run(); r.text = txt
                r.font.size = Pt(fs); r.font.color.rgb = RGBColor(*fg); r.font.bold = (ink > 0.30)
                _rp2._setfont(r)
        else:
            s.shapes.add_picture(g, Cm(0), Cm(0), Cm(21.0), Cm(29.7))
    bio = io.BytesIO(); prs.save(bio)
    return bio.getvalue()


def remodel(old_bytes, new_bytes, client='고객', base_date=''):
    o = read_sheet(old_bytes)
    n = read_sheet(new_bytes)
    c = compare(o, n)
    return c, None


# ─────────────────────────── 비교 엑셀 ───────────────────────────
def build_xlsx(cmp_, client='고객', base_date=''):
    """비교 엑셀(bytes). ★지점장 지시 「최종비교엑셀·리포트·보장분석지」의 첫째.
       수정을 생각하면 엑셀 — 값 확인·보정용이라 <b>수식이 아니라 실측값</b>을 쓴다."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    # ★★★★★v471 제76조 (지점장 2026.08.17 「색들만 더 눈에 띄게 하자」)
    #   기존 색은 흰 배경에서 거의 안 보였다 — 표 머리 #F2F5F9, 신규 #EAF6EF.
    #   ★진하게 올린다. 인쇄해도 구분이 남아야 한다.
    NAVY = 'FF0B2340'; GOLD = 'FFC5A052'
    HDR = PatternFill('solid', fgColor=NAVY)
    SUB = PatternFill('solid', fgColor='FFD6E2F0')    # 표 머리 — 옅은 회청 → 또렷한 하늘남
    NEWF = PatternFill('solid', fgColor='FFC8EBD8')   # 신규 — 연그린 → 진한 민트
    UPF = PatternFill('solid', fgColor='FFDDEBFB')    # 보장 증가 — 하늘
    DNF = PatternFill('solid', fgColor='FFFDE2C8')    # 보장 감소 — 주황
    DELF = PatternFill('solid', fgColor='FFF7CFCF')   # 삭제 — 분홍
    GRPF = PatternFill('solid', fgColor='FFE9EEF5')   # 구분(A열 그룹) — 회청
    # ★글자 포인트 (지점장 지시 2026.08.15) — 기본 11pt는 A4 표에서 크다
    _FS = 9
    W = Font(color='FFFFFFFF', bold=True, name='맑은 고딕', size=10)
    B = Font(bold=True, name='맑은 고딕', size=_FS)
    N = Font(bold=True, name='맑은 고딕', size=_FS)   # ★전부 진하게
    G = Font(bold=True, color='FF1F7A4D', name='맑은 고딕', size=_FS)
    R = Font(bold=True, color='FFC0444C', name='맑은 고딕', size=_FS)
    thin = Side(style='thin', color='FFD9DEE6'); BD = Border(thin, thin, thin, thin)
    med = Side(style='medium', color='FF0B2340')          # ★구분 경계선
    C = Alignment('center', 'center')
    CV = Alignment('center', 'center', wrap_text=True)    # 구분 셀(병합) 세로 가운데

    wb = openpyxl.Workbook()
    ws = wb.active; ws.title = 'MAKEONE LIFE PLAN'   # ★시트명 (지점장 지시 2026.08.15)
    ws.column_dimensions['A'].width = 1.6
    for col, w in zip('BCDEFG', [12, 30, 13, 13, 11, 9]):
        ws.column_dimensions[col].width = w

    def band(r, txt, span=5):   # ★B~G — 모든 표의 오른쪽 끝을 맞춘다
        ws.cell(r, 2, txt).font = W
        for c in range(2, 3 + span):
            ws.cell(r, c).fill = HDR; ws.cell(r, c).border = BD
        ws.merge_cells(start_row=r, start_column=2, end_row=r, end_column=2 + span)
        ws.cell(r, 2).alignment = C          # ★섹션 바도 가운데(지점장 「다 가운데」)

    def _pfoot(r0):
        """페이지 꼬리 = 골드 줄 + 네이비 바 (1쪽 하단과 같은 모양)"""
        for c in range(1, 8): ws.cell(r0, c).fill = GOLDF
        ws.row_dimensions[r0].height = 9
        for c in range(1, 8): ws.cell(r0 + 1, c).fill = HDR
        ws.cell(r0 + 1, 2, 'MAKEONE  보장분석 자동화').font = Font(bold=True, size=9, color='FFE6C878', name='맑은 고딕')
        ws.cell(r0 + 1, 7, f'{client} 고객님').font = Font(size=9, color='FFFFFFFF', name='맑은 고딕')
        ws.cell(r0 + 1, 7).alignment = Alignment('right', 'center')
        ws.row_dimensions[r0 + 1].height = 20
        return r0 + 2

    def _phead(r0):
        """페이지 머리 = 네이비 바(제목) + 골드 줄 (1쪽 상단과 같은 모양)"""
        for c in range(1, 8): ws.cell(r0, c).fill = HDR
        ws.cell(r0, 2, f'MAKEONE LIFE PLAN  —  {client} 고객님').font = \
            Font(bold=True, size=13, color='FFFFFFFF', name='맑은 고딕')
        ws.merge_cells(start_row=r0, start_column=2, end_row=r0, end_column=4)
        ws.cell(r0, 2).alignment = Alignment('left', 'center')
        ws.cell(r0, 6, f'제안 기준일 {base_date}').font = \
            Font(bold=True, size=9, color='FFE6C878', name='맑은 고딕')
        ws.merge_cells(start_row=r0, start_column=6, end_row=r0, end_column=7)
        ws.cell(r0, 6).alignment = Alignment('right', 'center')
        ws.row_dimensions[r0].height = 32
        for c in range(1, 8): ws.cell(r0 + 1, c).fill = GOLDF
        ws.row_dimensions[r0 + 1].height = 9
        return r0 + 2

    # ★★★★★v422q — 진단서처럼 <b>위·아래 포인트 바</b> (지점장 지시 2026.08.15).
    #   진단서 `.cvbar`(네이비 14mm + 골드 2.4mm) / `.cvfootbar`(네이비 6mm)를 엑셀 행으로 옮긴다.
    # ★★★★★v422s — 머리말을 <b>한 덩어리</b>로 (지점장 표시 2026.08.15).
    #   구 1~5행(바·골드·제목·기준일·빈줄)이 화면 위 5행을 먹고 표가 아래로 밀렸다.
    #   → 네이비 바 <b>한 줄 안</b>에 제목과 기준일을 같이 넣는다. 5행 → 2행.
    GOLDF = PatternFill('solid', fgColor=GOLD)
    for c in range(1, 8):
        ws.cell(1, c).fill = HDR
    # ★제목 (지점장 지시 2026.08.15 「리모델링 비교 대신 MAKEONE LIFE PLAN」)
    ws.cell(1, 2, f'MAKEONE LIFE PLAN  —  {client} 고객님').font = \
        Font(bold=True, size=13, color='FFFFFFFF', name='맑은 고딕')
    ws.merge_cells(start_row=1, start_column=2, end_row=1, end_column=4)
    ws.cell(1, 2).alignment = Alignment('left', 'center')
    ws.cell(1, 6, f'제안 기준일 {base_date}').font = \
        Font(bold=True, size=9, color='FFE6C878', name='맑은 고딕')
    ws.merge_cells(start_row=1, start_column=6, end_row=1, end_column=7)
    ws.cell(1, 6).alignment = Alignment('right', 'center')
    ws.row_dimensions[1].height = 32
    for c in range(1, 8):
        ws.cell(2, c).fill = GOLDF                     # 골드 포인트 줄
    ws.row_dimensions[2].height = 9                    # ★5는 모바일에서 안 보인다(실측 2026.08.15)

    # ★★★★★v422j (지점장 지적 2026.08.15 「엑셀에 전과후가없다」)
    #   계약 단위 <b>전 → 후</b> 표. 담보별 표만 있으면 「어느 보험이 빠지고 무엇이 들어왔나」가 안 보인다.
    band(4, '계약별 전 · 후')
    # ★v480 제89조 레이아웃 상수 — 이 값이 바뀌면 조문검사가 잡는다(제31조·제90조).
    # ★★★★★v484 제95조 (지점장 지적 2026.08.19 「한 칸씩 밀린다 그래서 또 2페이지가 된다」)
    #   인쇄 가능 높이는 769.7pt인데 예산을 <b>770pt로 꽉 채워</b> 잡았다 — 여유가 <b>0</b>이다.
    #   글꼴 렌더 반올림·프린터 드라이버 차이로 <b>한 줄이 툭 넘어가</b> 페이지가 하나 더 생긴다.
    #   ⇒ <b>한 줄(15pt) 여유</b>를 둔다. 페이지당 담보가 1개 줄지만 <b>절대 안 넘긴다</b>.
    #   ★한계값을 예산으로 쓰지 않는다. 예산은 한계보다 한 줄 작아야 한다.
    # ★★★★★v487 제97조 (지점장 실측 2026.08.19 「한 칸씩 밀린다 · 또 2페이지가 된다」)
    #   내 로컬 LibreOffice 변환은 3장인데 <b>지점장 화면에서는 6장</b>이었다.
    #   여유 1줄(15pt)로는 <b>뷰어별 행 높이 반올림</b>을 못 이긴다 — 한 줄이 밀리면
    #   그 줄과 꼬리가 <b>새 장</b>을 만들어 장수가 두 배가 된다.
    #   ⇒ 여유를 <b>3줄(45pt)</b>로 넓힌다. 꽉 채우는 것보다 <b>안 넘치는 것</b>이 먼저다.
    # ★★★★★v489 제97조 (지점장 실측 2026.08.19 「한 칸씩 밀린다 · 또 2페이지가 된다」)
    #   [실측] 지점장 뷰어에서 <b>6장</b>. 내 LibreOffice 변환은 3장 — <b>뷰어마다 달랐다</b>.
    #   [원인] 인쇄 가능 높이를 769.7pt로 잡았는데, 뷰어에 따라 <b>머리글·바닥글 여백
    #     (각 0.2in = 14.4pt)</b>을 본문에서 뺀다 → 실제 <b>740.9pt</b>.
    #     760pt는 769.7 안에는 들어가지만 <b>740.9는 넘는다</b> → 쪽마다 한 행씩 밀려 장수가 배로.
    #   [수정] 두 겹으로 막는다 — ㉠머리글·바닥글 여백을 <b>0</b>으로 못박고
    #     ㉡예산을 <b>735pt</b>로 낮춘다(769.7 대비 <b>2행 여유</b>).
    #   ★<b>「내 뷰어에서 3장」은 「모든 뷰어에서 3장」이 아니다.</b>
    _PAGE_PT = 735.0        # A4 인쇄 가능 769.7pt − 안전 여유 34.7pt(약 2행) ★v489
    _ROW_PT  = 15.0         # 담보 한 줄 기본 높이
    _FOOT_PT = 29.0         # 페이지 꼬리 = 골드 9pt + 네이비 20pt
    _SUM_PT_R = 15.0 + 5 * 15.0   # ★v485 요약표(밴드 1행 + 5행) 높이
    _LAYOUT89 = {'회사': 2, '상품명': (3, 4), '전': 5, '후': 6, '상태': 7}
    _PROD_W = 43            # ★v482 상품명 칸 표시폭 = C(30) + D(13)
    # ★★★★★v480 제89조 (지점장 지시 2026.08.18)
    #   「<b>엑셀 제일 처음 상품명을 한 칸 더 쓰고 상태를 제일 끝에 한 칸 쓰게 한 칸 줄여줘</b>」
    #   [구] B회사 · C상품명 · D전 · E후 · <b>F+G 상태(2칸 병합)</b>  → 상품명이 좁아 잘렸다
    #   [신] B회사 · <b>C+D 상품명(2칸)</b> · E전 · F후 · <b>G 상태(1칸)</b>
    #   ★열 배치를 바꾸면 <b>헤더·데이터·합계 세 곳</b>을 같이 바꾼다(제0조 전수 수정).
    for j, h in (2, '보험사'), (3, '상품명'), (5, '전 (기존)'), (6, '후 (변경 후)'), (7, '상태'):
        ws.cell(5, j, h).font = B
    for j in range(2, 8):
        ws.cell(5, j).fill = SUB; ws.cell(5, j).border = BD; ws.cell(5, j).alignment = C
    ws.merge_cells(start_row=5, start_column=3, end_row=5, end_column=4)   # ★상품명 2칸
    _r = 6
    _kk = {(c['company'], c['product']) for c in cmp_['keep']}
    _pk = {(c['company'], c['product']) for c in cmp_['prop']}
    for c in cmp_['old']['contracts'] + cmp_['prop']:
        key = (c['company'], c['product'])
        # ★★★★★v539 제125조 (지점장 실측 2026.08.22 「기존계약에 금액이 줄었는데 동일하게 나온다.
        #   감액으로 나와야 한다」 · 정답지 대조 — 119,010→43,369이 「유지 119,010」으로 나갔다)
        #   [원인] 유지 계약을 <b>before = after = 기존 보험료</b>로 박아
        #     <b>최종 엑셀의 보험료를 아예 읽지 않았다</b>. 「감액」 상태가 구조적으로 나올 수 없었다.
        #   ⇒ 유지 계약도 <b>최종 엑셀의 보험료</b>를 after로 쓴다. 줄면 <b>감액</b>, 늘면 <b>증액</b>.
        _newp = None
        for _nc in cmp_['new']['contracts']:
            if (_nc.get('company'), _nc.get('product')) == key:
                _newp = _nc.get('premium'); break
        if key in _pk:
            before, after, tag, fn = 0, c['premium'], '신규', G
        elif key in _kk or not cmp_['kill']:
            before = c['premium']
            after  = _newp if _newp is not None else c['premium']
            if   round(after) < round(before): tag, fn = '감액', R
            elif round(after) > round(before): tag, fn = '증액', G
            else:                              tag, fn = '유지', N
        else:
            before, after, tag, fn = c['premium'], 0, '삭제', R
        if key not in _pk and any((k['company'], k['product']) == key for k in cmp_['kill']):
            before, after, tag, fn = c['premium'], 0, '삭제', R
        # ★★★★★v422r — 신규 강조 (지점장 지시 2026.08.15).
        #   리모델링에서 <b>새로 들어온 것</b>이 결론이다. 같은 글꼴로 묻히면 안 된다.
        _new = (tag == '신규')
        ws.cell(_r, 2, c['company']).font = (B if _new else N)
        ws.cell(_r, 3, c['product']).font = (B if _new else N)   # ★v480 C+D 2칸
        ws.cell(_r, 5, round(before)).font = N
        ws.cell(_r, 6, round(after)).font = (G if _new else B)
        ws.cell(_r, 7, tag).font = fn                            # ★v480 상태 = 끝 1칸
        for cc in (5, 6): ws.cell(_r, cc).number_format = '#,##0"원"'
        # ★v539 제125조 2항 (지점장 「삭제나 감액에 색이 표시되어야 한다」)
        #   구 코드는 <b>신규만</b> 줄 색을 깔았다. 감액·증액·삭제는 글자색뿐이라 눈에 안 들어왔다.
        _cfill = (NEWF if tag == '신규' else DNF if tag == '감액'
                  else UPF if tag == '증액' else DELF if tag == '삭제' else None)
        for cc in range(2, 8):
            ws.cell(_r, cc).border = BD
            if _cfill: ws.cell(_r, cc).fill = _cfill
        ws.merge_cells(start_row=_r, start_column=3, end_row=_r, end_column=4)
        for cc in range(2, 8): ws.cell(_r, cc).alignment = C
        ws.cell(_r, 3).alignment = Alignment('center', 'center', wrap_text=True)  # ★잘림 방지
        # ★★★★★v482 제92조 (지점장 산출물 실측 2026.08.19) — <b>상품명이 두 줄이면 행도 두 줄이다</b>.
        #   제91조가 모든 행을 15pt로 못박은 뒤, <b>칸보다 긴 상품명</b>은 wrap_text로 두 줄이 되는데
        #   행 높이가 15pt에 묶여 <b>아랫줄이 잘렸다</b>(실측 KB `…(24.05)_1형_연만기` 표시폭 63 / 칸 43).
        #   → 표시폭(한글 2 · 영숫자 1)으로 줄 수를 세어 <b>필요한 만큼 높이를 준다</b>.
        #   ★높이 예산(제80조)은 실제 행 높이를 읽으므로 이 값이 <b>페이지 계산에 그대로 반영</b>된다.
        _disp = sum(2 if ord(_ch) > 0x1100 else 1 for _ch in str(c['product'] or ''))
        _ln = max(1, -(-_disp // _PROD_W))
        ws.row_dimensions[_r].height = _ROW_PT * _ln
        _r += 1
    ws.cell(_r, 2, '합계').font = B
    ws.cell(_r, 5, round(cmp_['prem_old'])).font = B                # ★v480
    ws.cell(_r, 6, round(cmp_['prem_new'])).font = B
    for cc in (5, 6): ws.cell(_r, cc).number_format = '#,##0"원"'
    for cc in range(2, 8): ws.cell(_r, cc).border = BD; ws.cell(_r, cc).fill = SUB
    ws.merge_cells(start_row=_r, start_column=3, end_row=_r, end_column=4)
    for cc in range(2, 8): ws.cell(_r, cc).alignment = C
    _r += 2

    # ★v471 — 상단 KPI 6칸은 뺐다. 지점장 「그냥 이거 쓰고 색들만 더 눈에 띄게」(2026.08.17).
    band(_r, '보험료')
    rows = [('기존 보험료 합계', cmp_['prem_old']), ('삭제 후 보험료(유지 계약)', cmp_['prem_keep']),
            ('제안 보험료(신규 계약)', cmp_['prem_prop']), ('최종 리포트 금액', cmp_['prem_new']),
            (('월 절감금액' if cmp_['save_m'] > 0 else '월 증가금액'), abs(cmp_['save_m'])),
            (('연 절감금액' if cmp_['save_m'] > 0 else '연 증가금액'), abs(cmp_['save_y']))]
    for i, (k, v) in enumerate(rows, start=_r + 1):
        ws.cell(i, 2, k).font = B; ws.cell(i, 5, round(v)).font = (G if ('절감' in k or '증가' in k) else N)
        ws.cell(i, 5).number_format = '#,##0"원"'
        for c in range(2, 8): ws.cell(i, c).border = BD
        ws.merge_cells(start_row=i, start_column=2, end_row=i, end_column=4)
        ws.merge_cells(start_row=i, start_column=5, end_row=i, end_column=7)
        ws.cell(i, 2).alignment = Alignment('left', 'center', indent=1)   # ★앞이 비어 보인다
        ws.cell(i, 5).alignment = C
    _pr = _r + 7
    ws.cell(_pr, 2, '월 보험료 절감률' if cmp_['save_m'] > 0 else '월 보험료 증가율').font = B
    ws.cell(_pr, 5, abs(cmp_['save_pct']) / 100).font = G
    ws.cell(_pr, 5).number_format = '0.0%'
    for c in range(2, 8): ws.cell(_pr, c).border = BD
    ws.merge_cells(start_row=_pr, start_column=2, end_row=_pr, end_column=4)
    ws.merge_cells(start_row=_pr, start_column=5, end_row=_pr, end_column=7)
    ws.cell(_pr, 2).alignment = Alignment('left', 'center', indent=1)
    ws.cell(_pr, 5).alignment = C

    # ★★★★★v422v — 「2페이지니까 2페이지로 나눠야지」(지점장 2026.08.15).
    #   자동 축소에 맡기면 어디서 끊길지 모른다. <b>담보별 표 앞에서 명시적으로 끊는다</b>.
    #     1쪽 = 계약별 전·후 + 보험료   /   2쪽 = 담보별 전·후 101행 + 요약
    # ★★★★★v422x — 나눠지는 자리마다 <b>1쪽과 같은 틀</b>(지점장 지시 2026.08.15).
    from openpyxl.worksheet.pagebreak import Break
    # ★★★★★v422z — 1쪽 하단에 <b>사망 · 후유</b>까지 싣는다(지점장 지시 2026.08.15).
    #   1쪽이 25행뿐이라 아래가 비었다. 담보표를 1쪽에서 시작하고 <b>암 블록 앞</b>에서 끊는다.
    _pp = _pr + 1
    band(_pp, '담보별 전 · 후 (증감)')
    # ★★★★★v470 (지점장 확정 2026.08.17) — 상담 3열(검토 결과·설계사 설명·고객 결정)은
    #   <b>넣지 않는다</b>. 세로 A4에 6열이 정본이다. 시안을 그대로 베끼지 않는다.
    for j, h in enumerate(['구분', '담보', '전 (기존)', '후 (변경 후)', '증감', '변화'], 2):
        ws.cell(_pp + 1, j, h).font = B; ws.cell(_pp + 1, j).fill = SUB
        ws.cell(_pp + 1, j).border = BD; ws.cell(_pp + 1, j).alignment = C
    r = _pp + 2
    # ★★★★★v422p — 「변동 없음」은 담보별 표에서 <b>뺀다</b> (지점장 「A4에 넘친다」 2026.08.15).
    #   변하지 않은 담보는 <b>리모델링의 결과가 아니다</b>. 41행이 표의 87%를 먹어 A4를 넘겼다.
    #   개수는 아래 「요약」에 그대로 남는다 — 값이 사라지는 게 아니라 자리를 옮기는 것이다.
    # ★★★★★v422w — 담보표가 한 장을 넘는다(지점장 표시 2026.08.15). 수술 블록 앞에서 한 번 더 끊는다.
    #   행 번호가 아니라 <b>담보명으로</b> 자리를 찾는다(제11조 구조 가정 금지).
    # ★★★★★v422y — 페이지를 <b>균등하게</b> 나눈다(지점장 「1페이지가 너무 긴데」 2026.08.15).
    #   A4 세로 한 장에 들어가는 것은 대략 <b>46행</b>이다. 담보 101개를 한 덩이로 두면 67행짜리 쪽이 생긴다.
    #   → 쪽당 담보 <b>38개</b>로 끊는다. 개수가 달라져도 자동으로 필요한 만큼 쪽이 생긴다.
    # ★★★★★v468 제75조 (지점장 2026.08.17 「원래 5페이지여야 한다」)
    #   쪽당 38개면 4쪽이 된다(실측: 나누기 36·80·124 → 4쪽).
    #   쪽당 <b>30개</b>로 줄이면 담보 101개가 5쪽으로 갈라진다. 한 쪽이 덜 빽빽해 상담 중에 읽기 쉽다.
    # ★★★★★v473 제79조 (지점장 지적 2026.08.18 「페이지 반을 못 넘긴다 · 그전에는 가득찼었다」)
    #   [실측] 산출물 `이창재_remodel_compare.xlsx` — 1쪽 51행 <b>770pt(27.2cm)</b>로 가득 차는데
    #          2~4쪽은 30행 <b>460pt(16.2cm)</b>뿐이었다. A4 세로 인쇄 가능 높이의 <b>60%</b>다.
    #   [원인] 쪽당 담보 개수를 <b>숫자로 박아</b>(_PER=24) 두었다. 1쪽은 앞에 계약별 표가 있어
    #          24개면 딱 찼지만, 2쪽부터는 그 표가 없는데도 <b>같은 24개</b>를 써서 절반이 비었다.
    #          검산: 머리41 + 밴드15 + 헤더15 + 담보 24×15 + 꼬리29 = 460pt — 실측치와 정확히 일치.
    #   [수정] 개수가 아니라 <b>높이 예산</b>으로 끊는다(제11조 구조 가정 금지).
    #          쪽마다 남은 높이를 실제 행 높이로 재서, 다음 담보행과 꼬리가 안 들어갈 때만 끊는다.
    #          ⇒ 1쪽 24개(종전 그대로) · 2쪽부터 44개. 쪽이 바뀌어도 자동으로 가득 찬다.

    def _pgpt(_a, _b):
        """_a행부터 _b-1행까지 실제 높이 합(pt). 미지정 행은 기본 15pt."""
        _t = 0.0
        for _x in range(_a, _b):
            _h = ws.row_dimensions[_x].height
            _t += float(_h) if _h else _ROW_PT
        return _t

    # ★★★★★v490 제100조 (지점장 지시 2026.08.19 「이 틀은 고정 · 그 안에 내용은 유동성」)
    #   [실측] 4쪽이 <b>머리 + 빈 행 35개 + 요약 + 꼬리</b>뿐이었다 — 담보가 하나도 없는 쪽.
    #   [지시] <b>머리·꼬리 틀은 페이지마다 고정</b>이고, 그 사이 <b>담보 수는 유동</b>이다.
    #   ⇒ 앞쪽을 꽉 채우고 마지막에 남는 것을 버리지 말고, <b>전체를 쪽수로 나눠 고르게</b> 싣는다.
    #   [계산] 쪽당 담보 예산 = 예산 − 머리(41) − 밴드·헤더(30) − 꼬리(29)
    #          마지막 쪽은 <b>요약(90) 자리를 빼고</b> 센다.
    _HEAD_PT2, _BANDHD_PT = 41.0, 30.0
    _per_cap = int((_PAGE_PT - _HEAD_PT2 - _BANDHD_PT - _FOOT_PT) // _ROW_PT)
    # ★★★★★v526 제110조 (지점장 지적 2026.08.21 「여백도 많다」 · 박주하 실측)
    #   구 v490은 <b>요약 자리를 빼지 않고</b> 쪽수를 나눴다. 그래서 마지막 쪽이 담보로 꽉 차고
    #   요약이 <b>또 한 쪽으로 밀려</b>, 그 쪽에 <b>요약만 남고 530pt(35행)가 비었다</b>(실측 로그
    #   `[v478 하단표] … 남은 530pt`). 앞쪽도 꼬리 고정 패딩 때문에 17행씩 비었다.
    #   ⇒ <b>마지막 쪽은 요약 자리를 뺀 한도</b>로 세고, 그 한도를 배분에 <b>미리</b> 반영한다.
    #     제100조(틀 고정 · 내용 유동)를 「요약도 틀」로 넓힌 것이다.
    #   ★요약 배치 판정(제87조)은 `_pgpt(_ptop, r+1)` — <b>아직 안 쓴 행 한 줄</b>을 더 세고,
    #     거기에 <b>안전 여유 한 줄</b>(v484 제95조)을 또 본다. 실측 3쪽 596pt인데 판정은 611pt.
    #     그 뒤 v484 제95조 판정(`_need`)은 <b>또 두 줄</b>(`_ROW_PT*2`)을 뺀다.
    #     ⇒ 실제로 통과하려면 한도에서 <b>세 줄</b>을 빼야 한다.
    #     실측 이력: 2줄 뺀 [27,38,36]·[27,39,35]·[27,40,34] 모두 요약이 밀렸다.
    _last_cap = int((_PAGE_PT - _HEAD_PT2 - _BANDHD_PT - _SUM_PT_R - _FOOT_PT - _ROW_PT * 3) // _ROW_PT)
    _n_all   = len(cmp_['all'])
    _first_cap = int((_PAGE_PT - _pgpt(1, r) - _FOOT_PT) // _ROW_PT)   # 1쪽은 앞 표들이 있다
    # ★★★★★v527 제100조 2항 (지점장 지시 2026.08.21 「마지막 2개에 일배책이 따로 분리다 ·
    #   공백이 펑 비어서 2페이지의 낭비다 · 늘 공백없이 보기좋게」)
    #   [실측 박주하] 계산은 「3쪽 · 배분 [25,43,33]」이라 찍고 실제는 <b>[25,42,33,1]</b>이었다.
    #   중간 몫 43이 <b>높이 한도 _per_cap=42를 넘어</b> 2쪽에서 높이 컷이 먼저 걸렸고,
    #   밀린 1개가 <b>4쪽에 홀로</b> 남아 그 쪽이 통째로 비었다(일상배상책임 1행).
    #   ⇒ ①쪽수는 <b>쪽별 실제 한도의 합</b>으로 구한다(구 while 루프는 마지막 쪽 한도를
    #     중간 쪽에도 섞어 써서 3쪽이라 오판했다) ②몫은 <b>한도를 절대 넘지 않게</b> 라운드로빈으로
    #     고르게 채운다. 계산과 실제가 어긋나면 그것이 결함이다(제91조).
    def _caps_for(_k):
        if _k <= 1: return [min(_first_cap, _last_cap)]
        return [_first_cap] + [_per_cap] * (_k - 2) + [_last_cap]
    _pages = 1
    while sum(_caps_for(_pages)) < _n_all and _pages < 60:
        _pages += 1
    _caps  = _caps_for(_pages)
    # ★v527 2항 (지점장 지적 2026.08.21 「비어있다」 · 화면 실측 79~95·126~143·174~182 빈 행)
    #   균등 배분은 <b>쪽마다 16~17행씩 빈 행</b>을 만든다(꼬리 바닥 고정 패딩, 제101조).
    #   ⇒ <b>최소 쪽수에서 앞쪽부터 한도까지 꽉 채운다.</b> 빈 행이 생기는 쪽은 마지막 하나뿐이다.
    # ★★★★★v527 제100조 2항 (지점장 지시 2026.08.21 「페이지는 유동성이라고 계속 말했다」)
    #   쪽수를 아끼려고 앞쪽에 몰면 <b>마지막 쪽이 담보 1~6개로 텅 빈다</b>(실측 4쪽 33행 공백).
    #   ⇒ <b>쪽수는 결과일 뿐이다.</b> 전체를 쪽수로 <b>고르게</b> 나누고, 각 쪽은 아래 3항
    #     (행 높이 확장)으로 <b>꽉 채운다</b>. 한도는 절대 넘지 않는다(넘으면 높이 컷이 먼저 걸린다).
    # ★★★★★v536 제100조 4항 (지점장 지시 2026.08.21
    #   「비교엑셀은 <b>절대 공백 주지마라</b> · 페이지는 줄여도 되니까」)
    #   [실측 101담보] 두 방식을 돌려 <b>빈 행 총수</b>를 셌다.
    #     ① 앞쪽 꽉 채우기 [25,42,28,6] → 빈행 [5,5,5,33] = <b>48</b>
    #     ② 고르게 나누기  [25,26,25,25] → 빈행 [5,5,7,10] = <b>27</b>
    #   ⇒ <b>②가 공백이 적다.</b> 쪽마다 조금씩 남기는 편이, 마지막 쪽에 몰아 <b>한 쪽을 통째로
    #     비우는 것</b>보다 낫다. 쪽수는 4쪽이 <b>물리적 최소</b>다(25+42+33=100 < 101 → 3쪽 불가).
    #   ★공백 총수가 판단 기준이다. 배분 방식을 바꿀 때는 <b>반드시 빈 행을 세어 비교</b>한다.
    # ★★★★★v538 제124조 5항 (지점장 지적 2026.08.21 「왜 칸들의 세로길이가 다 다르냐?」)
    #   구 v536 균등 배분은 쪽마다 <b>남는 높이가 달라</b> 행 높이가 15.1 / 24 / 24 / 20으로 제각각이었다.
    #   ⇒ 담보 수를 <b>쪽별 가용 높이(한도)에 비례</b>해 나눈다. 그러면 (가용높이 ÷ 담보수)가
    #     모든 쪽에서 같아져 <b>행 높이가 저절로 통일</b>되고, 공백도 0으로 유지된다.
    #   실측 101담보 · 한도 [25,42,42,33] → 배분 [18,30,30,23] → 높이 20.8 / 21.0 / 21.0 / 21.5
    #   ★1쪽은 위에 계약표·보험료표가 있어 <b>한도 계산이 실제보다 조금 후하다</b>.
    #     실측 편차 3.4pt(17.8 vs 21.2)가 남아 1쪽 몫에 <b>보정계수 0.86</b>을 준다.
    _wt = [(_c * 0.86 if _i == 0 else _c) for _i, _c in enumerate(_caps)]
    _tot_cap = sum(_wt) or 1
    _quota = [max(1, int(round(_n_all * _w / _tot_cap))) for _w in _wt]
    for _ci in range(_pages):                     # 한도 초과 금지
        _quota[_ci] = min(_quota[_ci], _caps[_ci])
    _d = _n_all - sum(_quota)                     # 반올림 보정
    _ci = 0
    while _d != 0 and _ci < _pages * 4:
        _k = _ci % _pages
        if _d > 0 and _quota[_k] < _caps[_k]: _quota[_k] += 1; _d -= 1
        elif _d < 0 and _quota[_k] > 1:          _quota[_k] -= 1; _d += 1
        _ci += 1
    print('[v538 높이통일] 담보 %d개 · %d쪽 · 한도 %s → 배분 %s (합 %d)'
          % (_n_all, _pages, _caps, _quota, sum(_quota)))



    _ptop = 1
    _dstart = 1
    _brks = []
    _qi = 0
    for _one in [None]:
        _pg, _gs = None, None
        _heads = []
        _cnt, _first = 0, True
        for _i, (grp, nm, o, n, d, tag) in enumerate(cmp_['all']):
            # ★v469 — 예전에는 1쪽을 「사망·후유장애」에서 끊었다. 담보 전수를 실으면
            #   1쪽이 담보 9행뿐이라 <b>절반이 빈다</b>(실측). 이제는 높이로만 끊는다.
            _cut = (_pgpt(_ptop, r) + _ROW_PT + _FOOT_PT > _PAGE_PT)
            # ★★★★★v485(폐기) — 구 「요약 자리 예약 컷」. v489 제97조 3항으로 대체됐다.
            #   [실측] 담보표가 3쪽 640pt에서 끝났는데 요약표(90) + 꼬리(29)가 안 들어가
            #   <b>4쪽이 생기고 3쪽은 83%</b>가 됐다. 요약 한 덩어리 때문에 쪽이 하나 더 났다.
            #   → <b>마지막 쪽이 될 구간에서는 요약 자리를 미리 예약</b>한다.
            #     남은 담보가 이 쪽에 다 들어가지만 요약이 안 들어가면 <b>지금 끊는다</b>.
            # ★★★★★v489 폐기 — 「마지막 쪽이면 요약 자리를 미리 예약하고 없으면 끊는다」는
            #   <b>끊어도 다음 쪽에서 같은 상황이 반복</b>돼 7행짜리 쪽을 연달아 만들었다(실측 7장).
            #   ⇒ <b>담보는 높이로만 끊는다.</b> 요약이 안 들어가면 <b>요약만</b> 다음 쪽으로 보낸다(아래).
            # ★★★★★v489 — <b>요약 예약 컷은 새 쪽 첫 담보에서는 발동하지 않는다</b>.
            #   `_islast and not _sumfit`이 매 담보마다 다시 참이 되어 <b>끊고 또 끊었다</b>
            #   (실측: 7행짜리 쪽이 3개 연속 · 총 7장). <b>_cnt(이 쪽에 그린 담보 수)</b>가 0이면 안 끊는다.
            _qcut = (_qi < len(_quota) and _cnt >= _quota[_qi])   # ★v490 균등 배분
            if _i and _cnt and (_cut or _qcut):
                if _gs is not None and r - 1 > _gs:        # 병합이 페이지를 넘지 않게 끊는다
                    ws.merge_cells(start_row=_gs, start_column=2, end_row=r - 1, end_column=2)
                _gs, _pg = None, None
                _first, _cnt = False, 0
                _qi += 1                                   # ★v490 다음 쪽 몫으로
                # ★★★★★v491 제101조 (지점장 지시 2026.08.19 「계약리스트가 유동성이라 괜찮다.
                #   <b>위아래 틀만 고정하라</b>」) — 꼬리가 담보 끝에 바로 붙어 <b>페이지 중간에 떴다</b>
                #   (실측 2·3쪽 65.6% → 꼬리가 2/3 지점). 머리는 늘 맨 위인데 꼬리만 떠다녔다.
                #   ⇒ 남은 높이만큼 빈 행을 넣어 <b>꼬리를 페이지 맨 아래에 고정</b>한다.
                #     제87조(요약 바닥 정렬)를 <b>모든 쪽의 꼬리</b>로 넓힌 것이다.
                # ★★★★★v527 제100조 3항 (지점장 지적 2026.08.21 「비어있다」 · 화면 실측
                #   79~95·126~143 빈 행) — 구 v491은 남은 높이를 <b>빈 행</b>으로 메워
                #   화면에 <b>17~18줄짜리 공백</b>이 그대로 보였다.
                #   ⇒ 빈 행 대신 <b>그 쪽 담보 행의 높이를 고르게 늘려</b> 표가 페이지를 채우게 한다.
                #     꼬리는 종전대로 맨 아래(제101조). 상한 24pt — 그 이상은 빈 행으로 남긴다.
                _room = _PAGE_PT - _pgpt(_ptop, r) - _FOOT_PT
                _nrow = r - _dstart
                if _room > 0 and _nrow > 0:
                    _hh = min(24.0, _ROW_PT + _room / _nrow)
                    for _rr in range(_dstart, r): ws.row_dimensions[_rr].height = _hh
                # ★★★★★v537 제124조 3항 (지점장 지시 2026.08.21 「공백이란 없다」)
                #   구 v491은 남은 높이를 <b>빈 행</b>으로 메워 꼬리를 페이지 바닥에 고정했다.
                #   인쇄에는 안 보여도 <b>엑셀 화면에는 그대로 보인다</b>(실측 27행).
                #   ⇒ <b>빈 행을 넣지 않는다.</b> 남는 높이는 위에서 <b>행 높이</b>로 이미 채웠고,
                #     상한(24pt)에 걸려 남는 부분은 <b>그냥 짧은 쪽</b>으로 둔다.
                #     제101조(꼬리는 맨 아래)보다 <b>공백 0이 우선</b>이다.
                r = _pfoot(r)                                  # 앞쪽 꼬리
                _brks.append(r - 1)
                _ptop = r                                      # ★v473 새 쪽의 첫 행
                r = _phead(r)                                  # 다음 쪽 머리
                band(r, '담보별 전 · 후 (증감) — 이어서')
                for j2, h2 in enumerate(['구분', '담보', '전 (기존)', '후 (변경 후)', '증감', '변화'], 2):
                    ws.cell(r + 1, j2, h2).font = B; ws.cell(r + 1, j2).fill = SUB
                    ws.cell(r + 1, j2).border = BD; ws.cell(r + 1, j2).alignment = C
                r += 2
                _dstart = r
            _cnt += 1
            _new = (tag == '신규 추가')
            # ★구분이 바뀌면 앞 그룹을 <b>병합</b>하고 경계에 굵은 선을 긋는다
            _head = (grp != _pg)
            if _head:
                if _gs is not None and r - 1 > _gs:
                    ws.merge_cells(start_row=_gs, start_column=2, end_row=r - 1, end_column=2)
                _gs, _pg = r, grp
                _heads.append(r)
                ws.cell(r, 2, grp).font = Font(bold=True, name='맑은 고딕', size=_FS, color='FF0B2340')
                ws.cell(r, 2).fill = GRPF          # ★v471 구분 열도 색으로 갈라 보이게
                ws.cell(r, 2).alignment = CV
            ws.cell(r, 3, nm).font = (B if _new else N)
            ws.cell(r, 4, o).font = N; ws.cell(r, 5, n).font = (G if _new else B)
            ws.cell(r, 6, d).font = (G if d > 0 else (R if d < 0 else N))
            ws.cell(r, 7, tag).font = (G if d > 0 else (R if d < 0 else N))
            # ★★★★★v471 제76조 — 변화 유형마다 <b>줄 전체</b>에 색을 깐다.
            #   숫자만 보고는 무엇이 바뀌었는지 안 보인다. 색이 먼저 눈에 들어와야 한다.
            _fill = (NEWF if tag == '신규 추가' else UPF if tag == '보장 증가'
                     else DNF if tag == '보장 감소' else DELF if tag == '삭제' else None)
            if _fill:
                for cc in range(2, 8): ws.cell(r, cc).fill = _fill
            elif _head:
                ws.cell(r, 2).fill = GRPF
            for c in range(4, 7): ws.cell(r, c).number_format = '#,##0'
            # ★구분이 바뀌는 행은 위쪽 굵은 선 — 불린으로 판정한다(Side 객체 비교는 실패한다)
            for c in range(2, 8):
                ws.cell(r, c).border = (Border(left=thin, right=thin, top=med, bottom=thin) if _head else BD)
            for c in range(3, 8): ws.cell(r, c).alignment = C
            ws.cell(r, 3).alignment = Alignment('center', 'center', wrap_text=True)
            r += 1

    if _gs is not None and r - 1 > _gs:
        ws.merge_cells(start_row=_gs, start_column=2, end_row=r - 1, end_column=2)
    # ★병합이 테두리를 지운다 — 구분 경계선은 <b>병합을 다 끝낸 뒤</b> 다시 긋는다(실측 2026.08.15)
    for _h in _heads:
        for c in range(2, 8):
            ws.cell(_h, c).border = Border(left=thin, right=thin, top=med, bottom=thin)
    for _b in _brks:
        ws.row_breaks.append(Break(id=_b))
    # ★★★★★v478 제87조 — <b>요약(하단표)은 페이지 바닥에 붙인다</b>(지점장 지시 2026.08.18
    #   「계속 어중간하게 넘어간다 · 다 붙일순없어? · 위·아래 제목이랑 하단표도 넣어주고」).
    #   [실측 서은옥] 3쪽 87% — 담보표가 134행에서 끝나고 요약표가 그 바로 밑에 붙어
    #   <b>페이지 아래가 붕 떴다</b>. 표는 위에 몰리고 여백이 아래에 남아 어중간해 보인다.
    #   → 남은 높이만큼 <b>빈 행으로 밀어</b> 요약표 + 꼬리가 페이지 <b>바닥에 딱 붙게</b> 한다.
    #   요약표 = 밴드 15pt + 5행 75pt = 90pt · 꼬리 29pt → 바닥에서 119pt를 확보한다.
    _SUM_PT = 15.0 + 5 * _ROW_PT
    # ★★★★★v489 — 요약이 이 쪽에 안 들어가면 <b>요약만</b> 다음 쪽으로 보낸다.
    #   담보를 끊어서 자리를 만들지 않는다(그러면 7행짜리 쪽이 연달아 생긴다).
    if _pgpt(_ptop, r + 1) + _SUM_PT + _FOOT_PT + _ROW_PT > _PAGE_PT:
        r = _pfoot(r)
        ws.row_breaks.append(Break(id=r - 1))   # ★즉시 적용 — _brks 루프(730행)는 이미 지나갔다
        _ptop = r
        r = _phead(r)
        print('[v489 요약] 자리가 없어 요약을 다음 쪽으로 넘긴다')
    _used_now = _pgpt(_ptop, r + 1)
    #   ★꼬리는 `max_row + 2`부터라 요약표와 꼬리 사이에 <b>빈 행 1개</b>가 더 들어간다 —
    #     그 15pt까지 빼야 페이지를 넘지 않는다(v478 1차 실측 775pt = 101% 초과).
    #   ★v484 제95조 — 마지막 쪽도 <b>한 줄 여유</b>를 남긴다(`_ROW_PT` 하나 더).
    # ★v527 제100조 3항 — <b>마지막 쪽</b> 담보 행도 같은 방식으로 늘린다(요약·꼬리 자리는 남긴다).
    _room_l = _PAGE_PT - _pgpt(_ptop, r) - _SUM_PT - _FOOT_PT - _ROW_PT * 5
    _nrow_l = r - _dstart
    if _room_l > 0 and _nrow_l > 0:
        _hl = min(24.0, _ROW_PT + _room_l / _nrow_l)
        for _rr in range(_dstart, r): ws.row_dimensions[_rr].height = _hl
        _used_now = _pgpt(_ptop, r + 1)
    _need = _PAGE_PT - _SUM_PT - _FOOT_PT - _ROW_PT * 2 - _used_now
    # ★★★★★v484 제95조 — <b>요약이 안 들어가면 페이지를 넘긴다</b>.
    #   담보 루프의 끊기(cut)는 <b>마지막 쪽에서는 일어나지 않는다</b> — 담보가 남으면 계속 쌓이고
    #   그 뒤에 요약(90pt)+꼬리(29pt)가 붙어 <b>마지막 쪽만 예산을 넘긴다</b>(실측 760pt).
    #   ⇒ 자리가 모자라면 꼬리를 찍고 새 쪽을 열어 요약을 얹는다. <b>쪽이 하나 늘어도 안 넘긴다.</b>
    if _need < 0:
        r = _pfoot(r)
        ws.row_breaks.append(Break(id=r - 1))   # ★즉시 적용 — _brks 루프는 이미 지나갔다
        _ptop = r
        r = _phead(r)
        _used_now = _pgpt(_ptop, r + 1)
        _need = _PAGE_PT - _SUM_PT - _FOOT_PT - _ROW_PT * 2 - _used_now
        print('[v484 하단표] 요약 자리가 모자라 새 쪽으로 넘겼다')
    # ★v537 제124조 3항 — 마지막 쪽도 <b>빈 행으로 밀지 않는다</b>. 요약은 담보표 바로 아래.
    if _need >= _ROW_PT:
        print('[v537 공백0] 마지막 쪽 여유 %.0fpt — 빈 행을 넣지 않고 요약을 표 바로 아래에 붙인다' % _need)
    band(r + 1, '요약')
    for i, (k, v) in enumerate([('보장 증가 항목', len(cmp_['up'])), ('신규 추가 특약', len(cmp_['add'])),
                                ('보장 감소 항목', len(cmp_['down'])), ('삭제 특약', len(cmp_['delete'])),
                                ('변동 없는 항목', len(cmp_['same']))], start=r + 2):
        ws.cell(i, 2, k).font = B; ws.cell(i, 5, v).font = N
        for c in range(2, 8): ws.cell(i, c).border = BD
        ws.merge_cells(start_row=i, start_column=2, end_row=i, end_column=4)
        ws.merge_cells(start_row=i, start_column=5, end_row=i, end_column=7)
        ws.cell(i, 2).alignment = Alignment('left', 'center', indent=1)   # ★앞이 비어 보인다
        ws.cell(i, 5).alignment = C

    # ★★★★★v422o — A4 인쇄 설정 (지점장 지적 2026.08.15 「A4에 넘친다」)
    #   엑셀은 기본이 「인쇄 설정 없음」이다. 그대로 인쇄하면 폭이 잘려 두 장으로 흩어진다.
    #   폭은 <b>무조건 한 장</b>(fitToWidth=1), 세로는 자연히 넘기게 둔다(fitToHeight=0).
    # ★★★★★v481 제91조 (지점장 지적 2026.08.18 「페이지 뒤에 2줄 나온다 · 딱 떨어지게」)
    #   [원인] 높이 예산(제80조)은 미지정 행을 <b>15pt로 가정</b>해 계산했는데, 엑셀은 높이를
    #   지정하지 않은 행을 <b>폰트 기준으로 자동 계산</b>한다(맑은 고딕 9pt ≈ 15.75pt).
    #   실측 — 138행이 미지정. 15.75pt면 1쪽 50행이 <b>794pt</b>가 되어 769.7pt를 넘고
    #   <b>1~2줄이 다음 장으로 밀린다</b>. 내 계산은 760pt였다 — <b>가정과 실제가 달랐다</b>.
    #   [수정] <b>모든 행에 높이를 명시</b>한다. 계산에 쓴 값과 파일에 박힌 값을 같게 만든다.
    #   ★가정으로 계산했으면 그 가정을 <b>파일에 못박아</b> 실제와 일치시킨다.
    _fr = ws.max_row + 2
    for c in range(1, 8):
        ws.cell(_fr, c).fill = GOLDF                   # 하단 골드 포인트 줄
    ws.row_dimensions[_fr].height = 9
    for c in range(1, 8):
        ws.cell(_fr + 1, c).fill = HDR                 # 하단 네이비 바
    # ★v478 제87조 — 마지막 페이지 꼬리만 <b>A열</b>에 써서 글자가 안 보였다(A열 폭 1.6).
    #   다른 페이지 꼬리(`_pfoot`)와 <b>같은 B열</b>로 맞춘다.
    ws.cell(_fr + 1, 2, 'MAKEONE  보장분석 자동화').font = Font(bold=True, size=9, color='FFE6C878', name='맑은 고딕')
    ws.cell(_fr + 1, 7, f'{client} 고객님').font = Font(size=9, color='FFFFFFFF', name='맑은 고딕')
    ws.cell(_fr + 1, 7).alignment = Alignment('right', 'center')

    ws.sheet_format.defaultRowHeight = _ROW_PT
    _hfix = 0
    for _hr in range(1, ws.max_row + 1):
        if ws.row_dimensions[_hr].height is None:
            ws.row_dimensions[_hr].height = _ROW_PT; _hfix += 1
    print('[v481 행높이] 미지정 %d행에 %.0fpt 명시 — 계산과 실제를 일치시킨다' % (_hfix, _ROW_PT))

    ws.row_dimensions[_fr + 1].height = 20

    # ★병합·값 설정 뒤 A열 fill이 날아간다(실측 2026.08.15) → 마지막에 다시 칠한다
    for c in range(1, 8):
        ws.cell(1, c).fill = HDR
        ws.cell(2, c).fill = GOLDF
        ws.cell(_fr, c).fill = GOLDF
        ws.cell(_fr + 1, c).fill = HDR

    # ★★★★★v470 제75조 6항 (지점장 2026.08.17 「우린 세로야 / 고객검토·설계사검토 삭제다」)
    #   상담 3열을 뺐으므로 <b>세로로 되돌린다</b>. 우리 문서는 세로가 정본이다.
    ws.page_setup.orientation = 'portrait'
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0   # ★세로는 자동 — 나눔이 정한다(지점장 「1페이지가 너무 긴데」)
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    ws.page_margins.left = ws.page_margins.right = 0.4
    ws.page_margins.top = ws.page_margins.bottom = 0.5
    ws.page_margins.header = ws.page_margins.footer = 0    # ★v489 제97조 — 본문을 깎지 않게 0
    ws.print_options.horizontalCentered = True
    ws.print_area = 'A1:G%d' % ws.max_row
    ws.oddFooter.center.text = '&P / &N'
    ws.oddFooter.center.size = 8

    bio = io.BytesIO(); wb.save(bio)
    return _shared_strings(bio.getvalue())


def _shared_strings(data):
    """★★★★★v423 — openpyxl은 글자를 `inlineStr`로 저장한다.
       <b>모바일 엑셀 앱 상당수가 이걸 못 연다</b>(지점장 「안열린다」 2026.08.15).
       시트 XML을 열어 sharedStrings로 바꾼다. Content_Types·rels 등록이 빠지면 파일이 거부된다."""
    import zipfile, re
    zin = zipfile.ZipFile(io.BytesIO(data))
    if 'xl/sharedStrings.xml' in zin.namelist():
        return data
    sheet = zin.read('xl/worksheets/sheet1.xml').decode('utf-8')
    strings, idx = [], {}

    def rp(m):
        a, t = m.group(1), m.group(2)
        if t not in idx:
            idx[t] = len(strings); strings.append(t)
        return '<c%s t="s"><v>%d</v></c>' % (a.replace(' t="inlineStr"', ''), idx[t])

    new = re.sub(r'<c([^>]*?)\s+t="inlineStr"[^>]*><is><t[^>]*>(.*?)</t></is></c>',
                 rp, sheet, flags=re.S)
    new = re.sub(r'<c([^>]*?)\s+t="inlineStr"[^>]*><is/></c>', r'<c\1/>', new)
    if not strings:
        return data
    ss = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
          '<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
          'count="%d" uniqueCount="%d">%s</sst>'
          % (len(strings), len(strings),
             ''.join('<si><t xml:space="preserve">%s</t></si>' % x for x in strings)))
    ct = zin.read('[Content_Types].xml').decode('utf-8')
    if 'sharedStrings' not in ct:
        ct = ct.replace('</Types>', '<Override PartName="/xl/sharedStrings.xml" ContentType='
                        '"application/vnd.openxmlformats-officedocument.spreadsheetml.'
                        'sharedStrings+xml"/></Types>')
    rl = zin.read('xl/_rels/workbook.xml.rels').decode('utf-8')
    if 'sharedStrings' not in rl:
        rl = rl.replace('</Relationships>', '<Relationship Id="rIdSST" Type="http://schemas.'
                        'openxmlformats.org/officeDocument/2006/relationships/sharedStrings" '
                        'Target="sharedStrings.xml"/></Relationships>')
    bio = io.BytesIO()
    zo = zipfile.ZipFile(bio, 'w', zipfile.ZIP_DEFLATED)
    for nm in zin.namelist():
        zo.writestr(nm, new if nm == 'xl/worksheets/sheet1.xml'
                    else ct if nm == '[Content_Types].xml'
                    else rl if nm == 'xl/_rels/workbook.xml.rels' else zin.read(nm))
    zo.writestr('xl/sharedStrings.xml', ss)
    zo.close()
    return bio.getvalue()


def remodel_all(old_bytes, new_bytes, client='고객', base_date=''):
    """★지점장 지시 「최종비교엑셀 · 리포트 · 보장분석지」 — 앞의 둘을 만든다.
       보장분석지는 <b>최종 엑셀 자신</b>이므로 그대로 돌려준다(같은 파일이다)."""
    o = read_sheet(old_bytes); n = read_sheet(new_bytes)
    c = compare(o, n)
    _rp = build_report(c, client, base_date)
    return {'cmp': c,
            'xlsx': build_xlsx(c, client, base_date),
            'pdf': _rp[0], 'pngs': _rp[1],
            'pptx': build_report_pptx(_rp[1], _rp[0], client)}


def remodel_single(xlsx_bytes, client='고객', base_date='', totpg=3):
    """★★★★★v422h — 엑셀 <b>한 개</b>로 비교(지점장 확정 2026.08.15 「엑셀은 1번만 보면된다」).
       보유 계약 = 기존 / 보유+제안 = 최종. 산출은 remodel_all과 동일 3종."""
    o, n, has_prop = split_sheet(xlsx_bytes)
    c = compare(o, n)
    c['has_prop'] = has_prop
    _rp = build_report(c, client, base_date)
    return {'cmp': c,
            'xlsx': build_xlsx(c, client, base_date),
            'pdf': _rp[0], 'pngs': _rp[1],
            'pptx': build_report_pptx(_rp[1], _rp[0], client)}
