# -*- coding: utf-8 -*-
"""
fill_ppt_from_xlsx — 2026.06.11 정기철 검증분.
법칙 등식2: PPT 값 출처 = 완성 엑셀 합계열(K). let_engine 미경유.
규칙: tf.clear() 금지 / run-injection / 라벨 재타이핑 금지 / 모르면 빈칸.
검증 포인트(정기철에서 잡은 버그 3건 반영):
 1) 콜론 run 안에 후속 라벨('/ 요양병원 :')이 있으면 콜론 직후 삽입(간병인일당 오삽입 차단)
 2) 콜론 없는 1~5종 괄호칸은 run 직접 교체
 3) 표적치료 줄(값 2개)은 7.5pt 축소 — 줄바꿈→다빈치 겹침 차단
사용: python3 fill_ppt_from_xlsx.py 보장진단_고객.xlsx 고객명
"""
import sys
from pptx import Presentation
from pptx.util import Pt
from openpyxl import load_workbook

MASTER='MASTER_보장분석지_PPT_빈폼.pptx'

def won(v):
    v=int(v or 0)
    if v<=0: return ''
    eok,man=divmod(v,10000)
    if eok and man: return f'{eok}억{man:,}만'
    if eok: return f'{eok}억'
    return f'{man:,}만'

def _do(p,full,label,value,shrink=None):
    colon=full.find(':', full.find(label)+len(label))
    if colon<0: return False
    run_of=[]
    for ri,r in enumerate(p.runs): run_of+=[ri]*len(r.text)
    cr=run_of[colon] if colon<len(run_of) else len(p.runs)-1
    rt=p.runs[cr].text
    pos=sum(len(p.runs[i].text) for i in range(cr)); cin=colon-pos
    after=rt[cin+1:]
    if after.strip():                      # 같은 run에 다음 라벨 존재 → 콜론 직후 삽입
        p.runs[cr].text=rt[:cin+1]+' '+value+rt[cin+1:]
        if shrink: p.runs[cr].font.size=Pt(shrink)
        return True
    for ri in range(cr+1,len(p.runs)):
        p.runs[ri].text=value
        if shrink: p.runs[ri].font.size=Pt(shrink)
        return True
    p.runs[cr].text=rt+value; return True

def inject(shape,label,value,shrink=None):
    if not value: return False
    ps=[p for p in shape.text_frame.paragraphs if p.runs]
    for p in ps:
        full=''.join(r.text for r in p.runs)
        if full.lstrip().startswith(label) and _do(p,full,label,value,shrink): return True
    for p in ps:
        full=''.join(r.text for r in p.runs)
        if label in full and _do(p,full,label,value,shrink): return True
    return False

def setbox(shape,value):
    if not value: return False
    tf=shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text=value; return True
    return False

def paren(shape,text):                      # 1~5종 괄호칸(콜론 없음)
    for p in shape.text_frame.paragraphs:
        full=''.join(r.text for r in p.runs)
        if full.strip().startswith('(') and '/' in full:
            p.runs[0].text=f'({text})'
            for r in p.runs[1:]: r.text=''
            p.runs[0].font.size=Pt(8); return True
    return False

def run(xlsx,name,date=('','','')):
    wb=load_workbook(xlsx,data_only=True); ws=wb.active
    X=lambda r: ws.cell(r,11).value or 0    # K=합계열. 행번호는 허혈성진단 추가본(88행) 기준.
    prs=Presentation(MASTER); sl=prs.slides[0]
    by={sh.name:sh for sh in sl.shapes if sh.has_text_frame}
    done=0
    if setbox(by['TextBox 1'],name): done+=1
    for b,v in zip(['TextBox 16','TextBox 24','TextBox 25'],date):
        if v and setbox(by[b],v): done+=1
    M=[('TextBox 11','상해사망',won(X(10)),None),('TextBox 10','종신',won(X(6)),None),
       ('TextBox 8','상해3%',won(X(11)),None),
       ('TextBox 14','암진단비',won(X(16)),None),('TextBox 14','유사암',won(X(18)),None),
       ('TextBox 14','고액암',won(X(15)),None),('TextBox 14','항암치료',won(X(27)),None),
       ('TextBox 14','표적치료',won(X(19)),7.5),('TextBox 14','다빈치로봇수술비',won(X(24)),None),
       ('TextBox 49','산정특례',won(X(33)),None),('TextBox 49','혈전용해치료비',won(X(34)),None),
       ('TextBox 56','산정특례',won(X(39)),None),('TextBox 56','혈전용해치료비',won(X(43)),None),
       ('TextBox 17','질병수술',won(X(62)),None),('TextBox 17','뇌혈관수술',won(X(64)),None),
       ('TextBox 17','허혈성수술',won(X(65)),None),('TextBox 17','심장 수술',won(X(66)),None),
       ('TextBox 19','상해수술',won(X(55)),None),('TextBox 19','골절수술',won(X(59)),None),
       ('TextBox 19','화상수술',won(X(61)),None),
       ('TextBox 7','골절',won(X(76)+X(77)),None),('TextBox 7','화상',won(X(81)),None),
       ('TextBox 7','깁스',won(X(84)),None),('TextBox 5','일상배상책임',won(X(88)),None),
       ('TextBox 6','입원',won(X(85)),None),('TextBox 6','통원',won(X(86)),None),
       ('TextBox 22','상해일당',won(X(48)),None),
       ('TextBox 9','대인',won(X(70)),None),('TextBox 9','대물',won(X(71)),None),
       ('TextBox 9','합의금',won(X(72)),None),('TextBox 9','6주미만',won(X(73)),None),
       ('TextBox 9','변호사비',won(X(74)),None),('TextBox 9','자부상',won(X(75)),None)]
    for box,lab,val,sk in M:
        if box in by and inject(by[box],lab,val,sk): done+=1
    # 1~5종(종별 합산 슬래시는 엑셀 K57/K63에 이미 텍스트로 존재)
    q=ws.cell(63,11).value; s=ws.cell(57,11).value
    if isinstance(q,str) and paren(by['TextBox 17'],q): done+=1
    if isinstance(s,str) and paren(by['TextBox 19'],s): done+=1
    if setbox(by['TextBox 26'],won(X(28))): done+=1     # 뇌혈관진단 값박스
    for n,lab,v in [('TextBox 54','허혈성',won(X(44))),('TextBox 55','급성심근',won(X(41)))]:
        tf=by[n].text_frame
        if tf.paragraphs and tf.paragraphs[0].runs and v:
            r=tf.paragraphs[0].runs[0]; r.text=f'{lab} {v}'; r.font.size=Pt(9)
            tf.word_wrap=False; done+=1
    for nm in ['TextBox 1','TextBox 16','TextBox 24','TextBox 25','TextBox 26']:
        if nm in by:
            try: by[nm].text_frame.word_wrap=False
            except Exception: pass
    out=f'보장분석지_{name}.pptx'; prs.save(out)
    print('saved',out,'| 채운칸',done)

if __name__=='__main__':
    run(sys.argv[1], sys.argv[2] if len(sys.argv)>2 else '고객')
