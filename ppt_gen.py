# -*- coding: utf-8 -*-
"""채운 캐논 엑셀 → 보장요약 PPT(네이비+골드). gen_ppt(xlsx,out)."""
import json, openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
NAVY=RGBColor(0x0B,0x1F,0x3A); NAVY2=RGBColor(0x12,0x29,0x4B)
GOLD=RGBColor(0xC9,0xA1,0x4A); GOLD2=RGBColor(0xE6,0xC8,0x79); WHITE=RGBColor(0xEA,0xF0,0xF8); MUT=RGBColor(0x9F,0xB0,0xC8)
def _bg(s,c=NAVY):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c
def _box(s,x,y,w,h,t,sz,col,b=False,al=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=al; r=p.add_run(); r.text=t
    r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=col; r.font.name='맑은 고딕'; return tb
def gen_ppt(xlsx,out):
    R=json.load(open('canon_rows.json')); NR=R['namerow']; NC=R['NC']; LAST=R['last']
    ws=openpyxl.load_workbook(xlsx)['보장진단']
    cust=(ws.cell(1,1).value or '고객 보장진단').replace(' 보장진단','')
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    # 표지
    s=prs.slides.add_slide(blank); _bg(s,NAVY)
    s.shapes.add_shape(1,Inches(0),Inches(3.3),Inches(13.333),Inches(0.06)).fill.solid()
    _box(s,0,2.3,13.333,1,f'{cust} 님 보장분석',40,GOLD2,True,PP_ALIGN.CENTER)
    _box(s,0,3.5,13.333,0.6,'BARUM 최은혜 지점장 · 미래를 바르게 설계합니다',16,MUT,False,PP_ALIGN.CENTER)
    # 카테고리별 요약(합계>0만)
    from barum_dict import CANON
    cat_rows={}
    for cat,items in CANON.items():
        rows=[(nm,ws.cell(NR[nm],LAST).value) for nm in items if nm in NR]
        rows=[(nm,v) for nm,v in rows if isinstance(v,(int,float)) and v>0]
        if rows: cat_rows[cat]=rows
    cats=list(cat_rows.items())
    for k in range(0,len(cats),2):  # 슬라이드당 2카테고리
        s=prs.slides.add_slide(blank); _bg(s,NAVY)
        _box(s,0.6,0.4,12,0.7,f'{cust} 보장현황',24,GOLD2,True)
        x=0.6
        for cat,rows in cats[k:k+2]:
            _box(s,x,1.4,6,0.5,f'■ {cat}',18,GOLD,True)
            y=2.0
            for nm,v in rows[:11]:
                _box(s,x,y,4.2,0.4,nm,13,WHITE)
                _box(s,x+4.2,y,1.7,0.4,f'{int(v):,}만원',13,GOLD2,False,PP_ALIGN.RIGHT)
                y+=0.42
            x+=6.3
    prs.save(out); return {'slides':len(prs.slides._sldIdLst)}
if __name__=='__main__':
    print(gen_ppt('out_canon.xlsx','out_canon.pptx'))
