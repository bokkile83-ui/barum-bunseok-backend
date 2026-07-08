# -*- coding: utf-8 -*-
"""
report_pptx_shapes.py  —  BARUM 보장설명서 도형·텍스트 기반 10p 생성기 (B안)
정본 원칙: 값 출처 = 완성 엑셀 데이터셀(C ~ 끝열-1) 직접 합산. SUM 캐시 의존 금지.
"""
import sys, os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import openpyxl

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT = RGBColor(0xF2, 0xF3, 0xF5)
GREEN = RGBColor(0x1D, 0x7A, 0x46)
RED = RGBColor(0xC0, 0x39, 0x2B)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
FONT = "맑은 고딕"

W, H = Inches(13.333), Inches(7.5)


def _setfont(run, name):
    """latin + ea + cs 전부 지정. ea 누락 시 한글이 테마폰트로 치환됨."""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag[2:]}")
        if el is None:
            from pptx.oxml.ns import qn
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def box(slide, x, y, w, h, text, size=12, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, fill=None, line=None, anchor=MSO_ANCHOR.TOP, wrap=True):
    sh = slide.shapes.add_textbox(x, y, w, h)
    tf = sh.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.margin_top = tf.margin_bottom = Emu(22860)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; _setfont(r, FONT)
    return sh


def rect(slide, x, y, w, h, fill, line=None, radius=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def header(slide, name, page, subtitle):
    rect(slide, 0, 0, W, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), W, Inches(0.035), GOLD)
    box(slide, Inches(0.5), Inches(0.12), Inches(7), Inches(0.3),
        "MAKEONE · 보장분석 리포트", 10, True, GOLD)
    box(slide, Inches(0.5), Inches(0.42), Inches(8), Inches(0.55),
        f"{name} 고객님 보장 진단서", 22, True, WHITE)
    box(slide, Inches(11.2), Inches(0.18), Inches(1.7), Inches(0.5),
        str(page), 26, True, WHITE, PP_ALIGN.RIGHT)
    box(slide, Inches(10.2), Inches(0.7), Inches(2.7), Inches(0.3),
        subtitle, 9, False, GOLD, PP_ALIGN.RIGHT)


def footer(slide, name, page, total):
    rect(slide, 0, Inches(7.05), W, Inches(0.45), NAVY)
    box(slide, Inches(0.5), Inches(7.12), Inches(5), Inches(0.3),
        "MAKEONE 보장분석 자동화", 8, True, GOLD)
    box(slide, Inches(8), Inches(7.12), Inches(4.8), Inches(0.3),
        f"{name} 고객님 · {page} / {total}", 8, False, WHITE, PP_ALIGN.RIGHT)


def table(slide, x, y, cols, rows_data, widths, hdr_h=0.32, row_h=0.3, fs=9):
    """가벼운 표 — 도형 기반. rows_data: list[list[str|(str,RGBColor)]]"""
    cx = x
    for j, cw in enumerate(widths):
        rect(slide, cx, y, Inches(cw), Inches(hdr_h), NAVY)
        box(slide, cx, y, Inches(cw), Inches(hdr_h), cols[j], fs, True, WHITE,
            PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        cx += Inches(cw)
    yy = y + Inches(hdr_h)
    for i, row in enumerate(rows_data):
        cx = x
        bg = WHITE if i % 2 == 0 else LIGHT
        for j, cw in enumerate(widths):
            cell = row[j]
            txt, col = (cell if isinstance(cell, tuple) else (cell, BLACK))
            rect(slide, cx, yy, Inches(cw), Inches(row_h), bg,
                 line=RGBColor(0xDD, 0xDD, 0xDD))
            box(slide, cx, yy, Inches(cw), Inches(row_h), txt, fs, False, col,
                PP_ALIGN.CENTER if j else PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            cx += Inches(cw)
        yy += Inches(row_h)
    return yy


def _gauge_png(pct, path, px=440):
    """부분 채움 도넛 링 — PIL 픽셀 렌더. 투명 배경."""
    from PIL import Image, ImageDraw
    S = px * 4                                     # 4x 슈퍼샘플링 → 안티에일리어싱
    img = Image.new("RGBA", (S, S), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    col = ((0x1D, 0x7A, 0x46) if pct >= 70 else
           (0xC9, 0xA2, 0x27) if pct >= 40 else (0xC0, 0x39, 0x2B))
    ring = int(S * 0.085)
    pad = ring // 2 + 4
    bbox = [pad, pad, S - pad, S - pad]
    d.arc(bbox, 0, 360, fill=(0xE3, 0xE5, 0xE9, 255), width=ring)   # 트랙
    if pct > 0:
        d.arc(bbox, -90, -90 + 360 * pct / 100, fill=col + (255,), width=ring)
        if pct < 100:                                                # 진행 끝점 캡
            import math
            a = math.radians(-90 + 360 * pct / 100)
            r = (S - 2 * pad) / 2
            cx, cy = S / 2 + r * math.cos(a), S / 2 + r * math.sin(a)
            d.ellipse([cx - ring / 2, cy - ring / 2, cx + ring / 2, cy + ring / 2],
                      fill=col + (255,))
    img.resize((px, px), Image.LANCZOS).save(path)
    return path


def gauge(slide, cx, cy, pct, label, tmpdir):
    d = Inches(1.15)
    p = _gauge_png(pct, os.path.join(tmpdir, f"g{pct}_{abs(hash(label))%9999}.png"))
    slide.shapes.add_picture(p, cx - d / 2, cy - d / 2, d, d)
    box(slide, cx - d / 2, cy - Inches(0.16), d, Inches(0.32),
        f"{pct}%", 13, True, BLACK, PP_ALIGN.CENTER, wrap=False)
    box(slide, cx - Inches(0.9), cy + d / 2 + Inches(0.04), Inches(1.8), Inches(0.28),
        label, 9, True, BLACK, PP_ALIGN.CENTER, wrap=False)


# ────────────────────────────── 데이터 로드 ──────────────────────────────
def load(xlsx):
    wb = openpyxl.load_workbook(xlsx)
    ws = wb.active
    last = ws.max_column                     # 합계열
    ncon = last - 3                          # C..(last-1)
    name = str(ws.cell(1, 1).value).replace(" 보장진단", "").strip()

    contracts = []
    for c in range(3, last):
        raw = str(ws.cell(1, c).value or "").split("\n")
        comp = raw[0] if raw else ""
        prod = raw[1] if len(raw) > 1 else ""
        tag = raw[2] if len(raw) > 2 else ""
        prem = ws.cell(2, c).value or 0
        contracts.append(dict(comp=comp, prod=prod, tag=tag, prem=int(prem)))

    cov = {}
    for r in range(6, ws.max_row + 1):
        nm = ws.cell(r, 2).value
        if not nm:
            continue
        vals, slash = [], None
        for c in range(3, last):
            v = ws.cell(r, c).value
            if isinstance(v, str) and "/" in v:
                slash = v
            elif isinstance(v, (int, float)):
                vals.append(v)
        if slash:
            cov[str(nm).strip()] = slash
            continue
        # 합계열 수식이 정한 집계 규칙을 그대로 따른다 (SUM 일괄 적용 금지)
        f = str(ws.cell(r, last).value or "").upper()
        s = sum(vals) if vals else 0
        if "MAX(" in f:
            v = max(vals) if vals else 0
        elif "MIN(" in f:
            m = re.search(r",\s*(\d+)\s*\)\s*$", f)
            v = min(s, int(m.group(1))) if m else s
        elif f.startswith("=IF(SUM"):
            m = re.search(r">\s*0\s*,\s*(\d+)", f)
            v = (int(m.group(1)) if m else s) if s > 0 else 0
        else:
            v = s
        cov[str(nm).strip()] = v
    return name, contracts, cov, ncon


def won(v):
    """만원 단위 → 한글 단위"""
    if isinstance(v, str):
        return v
    v = int(v)
    if v == 0:
        return ""
    if v >= 10000:
        eok, man = divmod(v, 10000)
        return f"{eok}억{man:,}만" if man else f"{eok}억"
    return f"{v:,}만"


# ────────────────────────────── 슬라이드 ──────────────────────────────
def build(xlsx, out, tmpdir='/tmp/_g'):
    os.makedirs(tmpdir, exist_ok=True)
    name, contracts, cov, ncon = load(xlsx)
    prem_sum = sum(c["prem"] for c in contracts)
    ren = sum(1 for c in contracts if "[갱신]" in c["tag"])
    non = len(contracts) - ren
    TOT = 10

    prs = Presentation(); prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    # ── 1. 표지
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, W, H, NAVY)
    box(s, Inches(1), Inches(2.1), Inches(8), Inches(0.5),
        "M A K E O N E", 16, True, GOLD)
    box(s, Inches(1), Inches(2.7), Inches(9), Inches(1.1),
        "보장 진단서", 44, True, WHITE)
    box(s, Inches(1), Inches(3.9), Inches(9), Inches(0.6),
        f"{name} 고객님", 20, False, WHITE)
    rect(s, Inches(1), Inches(4.9), Inches(3.4), Inches(0.02), GOLD)
    stats = [("보유 계약", f"{len(contracts)} 건"),
             ("월 납입보험료", f"{prem_sum:,} 원"),
             ("갱신 / 비갱신", f"{ren} / {non}")]
    for i, (k, v) in enumerate(stats):
        x = Inches(1 + i * 3.6)
        box(s, x, Inches(5.25), Inches(3.3), Inches(0.3), k, 10, False, GOLD)
        box(s, x, Inches(5.55), Inches(3.3), Inches(0.5), v, 18, True, WHITE)
    box(s, Inches(1), Inches(6.6), Inches(10), Inches(0.3),
        "MAKEONE · 보장분석 자동화 리포트", 9, False, GREY)

    # ── 2. 보장 현황
    s = prs.slides.add_slide(blank); header(s, name, 1, "보장 현황")
    cats = [
        ("사망·후유", ["일반사망", "상해사망", "상해후유3%", "상해후유80%"]),
        ("암", ["고액암", "일반암", "유사암(갑.기.경.제)", "암주요치료비", "항암방사선약물", "중입자치료비"]),
        ("뇌혈관", ["뇌혈관진단비", "뇌졸증진단비", "뇌출혈진단비", "외상성뇌출혈"]),
        ("심장", ["허혈성 진단비", "급성심근경색", "부정맥", "심부전"]),
        ("수술비", ["상해수술비", "중대한상해수술비", "심장수술비", "상해 종수술비(1-5종)"]),
        ("운전자", ["합의금", "변호사", "대인", "대물", "6주미만", "자부상"]),
        ("입원·일당", ["간병인", "간호통합병동", "상해중환자실", "상해수술일당", "질병일당", "상해일당"]),
        ("실손·일배책", ["일상배상책임", "입원", "통원", "약값"]),
        ("골절·화상", ["중증화상진단비", "골절(치아파절제외)", "5대골절진단비", "화상진단비"]),
        ("응급실·독감", ["응급실(응급)", "독감"]),
    ]
    y = Inches(1.45)
    for i, (cat, keys) in enumerate(cats):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col * 6.35)
        yy = y + Inches(row * 1.06)
        rect(s, x, yy, Inches(1.55), Inches(0.92), NAVY)
        box(s, x, yy, Inches(1.55), Inches(0.92), cat, 10, True, WHITE,
            PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        items = [f"{k.split('(')[0]} {won(cov.get(k, 0))}" for k in keys if cov.get(k, 0)]
        rect(s, x + Inches(1.55), yy, Inches(4.6), Inches(0.92), LIGHT,
             line=RGBColor(0xDD, 0xDD, 0xDD))
        txt = "   ".join(items[:3]) + ("\n" + "   ".join(items[3:6]) if len(items) > 3 else "")
        box(s, x + Inches(1.6), yy, Inches(4.5), Inches(0.92),
            txt if items else "미가입", 9, False,
            BLACK if items else RED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, name, 1, TOT - 1)

    # ── 3. AI 진단
    s = prs.slides.add_slide(blank); header(s, name, 2, "AI 진단")
    box(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.35),
        "갱신형 — 보험료 인상 가능", 13, True, RGBColor(0x18, 0x5F, 0xA5))
    box(s, Inches(6.9), Inches(1.4), Inches(6), Inches(0.35),
        "비갱신형 — 만기까지 고정", 13, True, RED)
    yr = yn = Inches(1.85)
    for c in contracts:
        renew = "[갱신]" in c["tag"]
        x = Inches(0.5) if renew else Inches(6.9)
        yy = yr if renew else yn
        rect(s, x, yy, Inches(5.9), Inches(0.62), LIGHT, line=RGBColor(0xDD, 0xDD, 0xDD))
        box(s, x + Inches(0.1), yy + Inches(0.04), Inches(4.3), Inches(0.55),
            f"{c['comp']}\n{c['prod'][:34]}", 8.5, False, BLACK)
        box(s, x + Inches(4.4), yy, Inches(1.4), Inches(0.62),
            f"{c['prem']:,}원", 10, True,
            RGBColor(0x18, 0x5F, 0xA5) if renew else BLACK,
            PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        if renew: yr = yy + Inches(0.72)
        else:     yn = yy + Inches(0.72)
    footer(s, name, 2, TOT - 1)

    # ── 4. 핵심 보장 분석 (CI 3-state)
    s = prs.slides.add_slide(blank); header(s, name, 3, "핵심 보장 분석")
    CI_KEYS = ("중대한 암", "중대한 뇌졸증", "중대한 급성심근", "중대한CI적용")
    ci_amt = sum(v for k, v in cov.items()
                 if k in CI_KEYS and isinstance(v, (int, float)))
    prods = " ".join(c["prod"] for c in contracts)
    if ci_amt > 0:
        state, tcol, title = "ci", GOLD, "CI(중대질병 선지급형) 계약"
        body = ("주계약 사망보험금에서 선지급되는 구조다. 중대한 암·뇌졸증·급성심근경색 "
                "진단 시 사망보험금의 일부가 먼저 지급되고, 잔여분이 사망보장으로 남는다.")
    elif any(k in prods for k in ("CI보험", "GI보험", "리빙케어")):
        state, tcol, title = "check", GOLD, "[확인] CI 판정 보류"
        body = ("상품명에 CI·GI·리빙케어 표기가 있으나 '중대한OO' 담보값이 추출되지 않았다. "
                "증권 원본 확인 필요. 일반 보험으로 단정하지 않는다.")
    else:
        state, tcol, title = "none", RGBColor(0x18, 0x5F, 0xA5), "CI보험 아님 · 비CI 계약"
        body = ("이 고객은 CI(중대질병 선지급형) 계약이 아니다. 암·뇌·심 진단비는 진단 즉시 "
                "정액 100% 지급되는 일반형이며, 사망보험금 차감·선지급·잔여 개념이 없다. "
                "진단금 전액이 치료비로 쓰이고 사망보장은 별도로 유지된다.")
    rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.5), LIGHT,
         line=RGBColor(0xDD, 0xDD, 0xDD))
    box(s, Inches(0.75), Inches(1.6), Inches(11.8), Inches(0.4), title, 15, True, tcol)
    box(s, Inches(0.75), Inches(2.05), Inches(11.8), Inches(0.8), body, 10.5, False, BLACK)

    box(s, Inches(0.5), Inches(3.2), Inches(6), Inches(0.35),
        "주요 치료비 정리  TREATMENT BENEFITS", 12, True, NAVY)
    tb = [("암주요치료비", "암주요치료비"), ("비급여주요치료비", "하이클래스(암)"),
          ("순환계주요치료비", "2대 주요치료비"), ("산정특례(뇌혈관)", "산정특례뇌혈관"),
          ("산정특례(심장)", "산정특례심장")]
    for i, (lab, key) in enumerate(tb):
        x = Inches(0.5 + i * 2.5)
        v = cov.get(key, 0)
        rect(s, x, Inches(3.65), Inches(2.3), Inches(1.05), WHITE,
             line=GOLD if v else RGBColor(0xDD, 0xDD, 0xDD))
        box(s, x, Inches(3.78), Inches(2.3), Inches(0.3), lab, 9, True, GREY, PP_ALIGN.CENTER)
        box(s, x, Inches(4.12), Inches(2.3), Inches(0.45),
            won(v) if v else "미가입", 14 if v else 11, True,
            NAVY if v else RED, PP_ALIGN.CENTER)

    box(s, Inches(0.5), Inches(5.1), Inches(6), Inches(0.35),
        "보장 진단 코멘트  SUMMARY", 12, True, NAVY)
    rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.15), LIGHT,
         line=RGBColor(0xDD, 0xDD, 0xDD))
    box(s, Inches(0.75), Inches(5.65), Inches(11.8), Inches(0.85),
        f"{name} 고객님은 보유 {len(contracts)}건 · 월 {prem_sum:,}원의 보장을 운용하고 있다.\n"
        f"갱신 {ren}건 · 비갱신 {non}건. 값은 완성 엑셀 끝열에서 직접 읽었다(등식2).",
        10.5, False, BLACK)
    footer(s, name, 3, TOT - 1)

    # ── 5. 부위별 충족률
    s = prs.slides.add_slide(blank); header(s, name, 4, "부위별 충족률")
    def amt(*keys): return sum(cov.get(k, 0) for k in keys if isinstance(cov.get(k, 0), (int, float)))
    def cnt(*keys): return sum(1 for k in keys if cov.get(k, 0))
    areas = [
        ("암", amt("고액암", "일반암", "유사암(갑.기.경.제)", "암주요치료비", "항암방사선약물", "중입자치료비", "표적항암치료비"), 10000, "만"),
        ("운전자", cnt("대인", "대물", "합의금", "6주미만", "변호사", "자부상"), 4, "개"),
        ("실손·배상", cnt("입원", "통원", "약값", "일상배상책임"), 3, "개"),
        ("수술비", cnt("상해수술비", "중대한상해수술비", "심장수술비", "골절수술비", "5대골절수술비", "창상봉합술", "질병수술비", "상해 종수술비(1-5종)"), 4, "개"),
        ("뇌혈관", amt("뇌혈관진단비", "뇌졸증진단비", "뇌출혈진단비", "외상성뇌출혈"), 5000, "만"),
        ("사망·후유", amt("일반사망", "상해사망", "질병사망(80세)", "상해후유3%", "상해후유80%", "질병후유80%"), 15000, "만"),
        ("골절·화상", cnt("골절(치아파절포함)", "골절(치아파절제외)", "5대골절진단비", "화상진단비", "중증화상진단비", "반깁스", "깁스진단비"), 3, "개"),
        ("심장", amt("허혈성 진단비", "급성심근경색", "협심증", "부정맥", "심부전"), 3000, "만"),
        ("입원·일당", cnt("질병일당", "상해일당", "상해수술일당", "간호통합병동", "상해중환자실", "간병인"), 3, "개"),
        ("응급실·독감", cnt("응급실(응급)", "독감"), 2, "개"),
    ]
    calc = [(a, min(100, round(v / g * 100)) if g else 0, v, g, u) for a, v, g, u in areas]
    for i, (a, p, *_ ) in enumerate(calc):
        gauge(s, Inches(1.35 + (i % 5) * 2.55), Inches(2.15 + (i // 5) * 1.85), p, a, tmpdir)
    box(s, Inches(0.5), Inches(5.55), Inches(6), Inches(0.3),
        "충족률 산정 근거   보유 ÷ 40대 권장", 11, True, NAVY)
    rows = [(a, f"{v:,}{u}", f"{g:,}{u}",
             (f"{p}%", GREEN if p >= 70 else (GOLD if p >= 40 else RED)))
            for a, p, v, g, u in calc[:5]]
    table(s, Inches(0.5), Inches(5.95), ["영역", "보유", "권장(40대)", "충족률"],
          rows, [2.0, 1.6, 1.6, 1.1], row_h=0.2, fs=8)
    rows2 = [(a, f"{v:,}{u}", f"{g:,}{u}",
              (f"{p}%", GREEN if p >= 70 else (GOLD if p >= 40 else RED)))
             for a, p, v, g, u in calc[5:]]
    table(s, Inches(6.9), Inches(5.95), ["영역", "보유", "권장(40대)", "충족률"],
          rows2, [2.0, 1.6, 1.6, 1.1], row_h=0.2, fs=8)
    footer(s, name, 4, TOT - 1)

    # ── 6. 담보별 보장범위
    s = prs.slides.add_slide(blank); header(s, name, 5, "담보별 보장범위")
    has_brain = cov.get("뇌혈관진단비", 0) or cov.get("뇌졸증진단비", 0) or cov.get("뇌출혈진단비", 0)
    has_heart = cov.get("허혈성 진단비", 0) or cov.get("급성심근경색", 0)
    O, X = "●", "○"
    box(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.32), "뇌 — 질병코드별 커버", 12, True, NAVY)
    brain = [
        ("뇌출혈  I60~62", O if has_brain else X, O, O),
        ("뇌졸증·뇌경색  I63·65·66", O if cov.get("뇌졸증진단비", 0) else X, O, O),
        ("기타 뇌혈관질환  I64·67~69", X, O, O),
        ("뇌동맥류·정맥류  I71·72", X, O, X),
        ("선천 뇌혈관기형  Q28.0~28.3", X, X, O),
        ("외상성 뇌출혈  S06", X, X, O),
    ]
    table(s, Inches(0.5), Inches(1.78), ["질병 (코드)", "뇌혈관진단비", "순환계", "산정특례"],
          brain, [3.0, 1.4, 1.0, 1.0], row_h=0.28, fs=8.5)
    box(s, Inches(0.5), Inches(4.0), Inches(6.2), Inches(0.9),
        "· 외상성 뇌출혈(S06) = 뇌혈관진단비 미보장 → 산정특례 축only\n"
        "· 산정특례 = 진단 기반 별개 담보축 · 지급조건은 회사·약관별 [확인]", 8.5, False, GREY)

    box(s, Inches(7.0), Inches(1.4), Inches(6), Inches(0.32), "심장 — 질병코드별 커버", 12, True, NAVY)
    heart = [
        ("급성심근경색  I21~23", O if cov.get("급성심근경색", 0) else X, O, O, O),
        ("협심증  I20", O if cov.get("허혈성 진단비", 0) else X, O, O, O),
        ("기타·만성 허혈  I24·25", O if cov.get("허혈성 진단비", 0) else X, O, O, O),
        ("심장판막  I05·I34~37", X, O, O, O),
        ("심근·심내막 염증  I30~33·I40", X, O, O, O),
        ("부정맥  I49", X, O, O, O),
        ("심부전  I50", X, O, O, O),
        ("심근병증  I42~45", X, O, O, O),
    ]
    table(s, Inches(7.0), Inches(1.78), ["질병 (코드)", "허혈성진단비", "심장(특정)", "순환계", "산정특례"],
          heart, [2.5, 1.2, 1.0, 0.85, 0.85], row_h=0.28, fs=8.5)
    box(s, Inches(7.0), Inches(4.42), Inches(5.8), Inches(1.0),
        "· 빈맥(I47·48) ≠ 부정맥(I49) — 전용 행 구분\n"
        "· 순환계 = 2대 + 동맥류·정맥류 등 순환기 전반. 세부 대상코드 [확인]\n"
        "● 보장   ○ 미보장", 8.5, False, GREY)
    footer(s, name, 5, TOT - 1)

    # ── 7. 치료비 변천사
    s = prs.slides.add_slide(blank); header(s, name, 6, "주요치료비 변천사")
    box(s, Inches(0.5), Inches(1.35), Inches(12), Inches(0.35),
        "3대 주요치료비 담보의 변천사   ①비례형(구간) → ②정액형 → ③비급여 → 생활비", 12, True, NAVY)
    stages = [
        ("1  비례형(구간형)", "치료비 지출액이 구간 하한 도달 시 그 구간 정액 지급.\n2024.11 단종 — 보유자 해지 금지."),
        ("2  정액형", "치료 사실만으로 약정금 정액 지급.\n면책 90일 · 가입금액 100만원부터."),
        ("3  비급여(하이클래스)", "급여 전액본인부담금 + 비급여 치료비 정액.\n2천만×10년=2억 · 25.08~ 생활비형 추가."),
        ("4  왜 필요한가", "5세대 실손: 비중증 비급여 자기부담 50%·한도 1,000만.\n표적·면역 고가 약값엔 부족 → 주요치료비가 메움."),
    ]
    for i, (t, d) in enumerate(stages):
        x = Inches(0.5 + (i % 2) * 6.4); yy = Inches(1.85 + (i // 2) * 1.35)
        rect(s, x, yy, Inches(6.05), Inches(1.15), LIGHT, line=RGBColor(0xDD, 0xDD, 0xDD))
        box(s, x + Inches(0.12), yy + Inches(0.06), Inches(5.8), Inches(0.3), t, 11, True, GOLD)
        box(s, x + Inches(0.12), yy + Inches(0.38), Inches(5.8), Inches(0.7), d, 9, False, BLACK)
    box(s, Inches(0.5), Inches(4.75), Inches(12), Inches(0.32),
        "수술 건당 총진료비(평균) · 심평원·건보공단", 11, True, NAVY)
    surg = [("심장수술(개흉)", "2,832만원"), ("관상동맥우회술(CABG)", "2,690만원"),
            ("뇌동맥류 코일색전술", "1,100~1,600만원"), ("심장 스텐트(다발혈관)", "1,000~1,400만원"),
            ("뇌기저부수술", "1,475만원")]
    for i, (k, v) in enumerate(surg):
        x = Inches(0.5 + i * 2.5)
        rect(s, x, Inches(5.15), Inches(2.3), Inches(0.85), WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        box(s, x, Inches(5.24), Inches(2.3), Inches(0.32), k, 8.5, False, GREY, PP_ALIGN.CENTER)
        box(s, x, Inches(5.56), Inches(2.3), Inches(0.35), v, 12, True, NAVY, PP_ALIGN.CENTER)
    box(s, Inches(0.5), Inches(6.15), Inches(12), Inches(0.6),
        "재발률 — 뇌경색 1년 10%·5년 20~30% / 급성심근경색 1년 30%·3년 50% / 스텐트 5년 재협착 15%",
        9.5, False, GREY)
    footer(s, name, 6, TOT - 1)

    # ── 8. 상담 워크시트
    s = prs.slides.add_slide(blank); header(s, name, 7, "상담 워크시트")
    box(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.32), "암 주요치료비", 13, True, GOLD)
    box(s, Inches(7.0), Inches(1.4), Inches(6), Inches(0.32), "뇌·심장 주요치료비", 13, True, GOLD)
    brain_lbl = " · ".join(
        f"{n} {won(cov.get(k,0))}" for n, k in
        [("뇌혈관", "뇌혈관진단비"), ("뇌졸증", "뇌졸증진단비"), ("허혈성", "허혈성 진단비"),
         ("급성심근", "급성심근경색")] if cov.get(k, 0)) or "미가입"
    cards = [
        (0.5, "암 진단비", "걸렸을 때 일시금 (기본)", won(cov.get("일반암", 0)) or "미가입"),
        (0.5, "암 주요치료비", "수술·방사선·약물 정액 (100만~)", won(cov.get("암주요치료비", 0)) or "미가입"),
        (0.5, "하이클래스 (비급여)", "표적·면역·중입자 전액본인 커버", won(cov.get("하이클래스(암)", 0)) or "미가입"),
        (7.0, "뇌·심 진단비", "뇌혈관·허혈성 일시금", brain_lbl),
        (7.0, "2대 주요치료비", "수술·혈전용해·중환자실 (100만~)", won(cov.get("2대 주요치료비", 0)) or "미가입"),
        (7.0, "산정특례 (뇌 · 심장)", "진단만으로 지급 · 개별 담보",
         (won(cov.get("산정특례뇌혈관", 0)) or "미가입") + " / " + (won(cov.get("산정특례심장", 0)) or "미가입")),
    ]
    yl = yr2 = Inches(1.82)
    for x0, t, sub, val in cards:
        x = Inches(x0); yy = yl if x0 == 0.5 else yr2
        ok = "미가입" not in val
        rect(s, x, yy, Inches(5.9), Inches(1.15), WHITE, line=GOLD if ok else RGBColor(0xDD, 0xDD, 0xDD))
        box(s, x + Inches(0.15), yy + Inches(0.08), Inches(3.5), Inches(0.32), t, 11, True, NAVY)
        box(s, x + Inches(0.15), yy + Inches(0.42), Inches(3.5), Inches(0.6), sub, 8.5, False, GREY)
        vs = 11.5 if len(val) <= 12 else (9.5 if len(val) <= 24 else 8.5)
        box(s, x + Inches(3.55), yy, Inches(2.25), Inches(1.15), val, vs, True,
            NAVY if ok else RED, PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        if x0 == 0.5: yl = yy + Inches(1.3)
        else:         yr2 = yy + Inches(1.3)
    rect(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.85), LIGHT, line=GOLD)
    box(s, Inches(0.7), Inches(5.98), Inches(11.9), Inches(0.6),
        "부족한 담보 → 보완 추천 :", 12, True, NAVY)
    footer(s, name, 7, TOT - 1)

    # ── 9·10. 심혈관 담보 분류
    SETS = [
        [("한화손해보험 4가지", ["심혈관Ⅰ = 협심증 I20 · 기타급성허혈 I24 · 만성허혈 I25 · 빈맥 I47 · 심방세동 I48 · 기타부정맥 I49 · 심부전 I50",
                            "심혈관Ⅱ = 급성심근경색 I21~23 · 인공소생 심장정지 I46.0",
                            "심근병증 = I42 (확장성·비후성·제한성 포함)",
                            "심혈관특정질환 = Ⅰ에서 기타부정맥(I49)만 제외"]),
         ("DB손해보험 4가지", ["특정Ⅰ = 협심증·허혈 I20·24·25 + 주요심장염증 I30~41",
                          "특정Ⅱ = 급성심근경색 I21~23 · I46.0",
                          "특정Ⅲ = 판막질환 · 빈맥 I47 · 심방세동 I48 · 심부전 I50",
                          "순환계 3대 = I46.0 · 부정맥 I47~49 · 심부전 I50"]),
         ("KB손해보험 5가지", ["특정Ⅰ = 협심증·허혈 I20·24·25 · 빈맥 I47·48 · 심부전 I50 (염증X·부정맥X)",
                          "특정Ⅱ = 급성심근경색 I21~23 · I46.0",
                          "심근병증 = I42 / 심장판막질환 = 판막·I33·I38",
                          "기타심장부정맥 = I49 (I47·48은 대상 아님)"]),
         ("현대해상 6가지", ["허혈성 심장질환 = 협심증 I20 · I24 · I25",
                        "특정허혈 심장질환 = 급성심근경색 I21~23",
                        "특정Ⅰ = 빈맥 I47 · 심방세동 I48 · 심부전 I50",
                        "특정Ⅱ = 급성심근경색 I21~23 · I46.0",
                        "주요 심장염증 = I30~41 / 특정2대 = 전도장애 I44·45 + I49",
                        "※ 현대 심근병증 담보 없음"])],
        [("NH농협손해보험 5가지", ["심혈관특정Ⅰ = 협심증·허혈 · 빈맥 · 기타부정맥 I49 · 심부전",
                            "특정Ⅰ(기타부정맥 제외) = Ⅰ에서 I49만 뺀 묶음",
                            "심근병증 = I42·I43 / 주요 심장염증질환 = I30~41",
                            "기타 심장부정맥 = I49"]),
         ("흥국화재 실제 약관", ["특정심혈관질환(기타부정맥제외) = 협심증·허혈·빈맥·심부전 (급성심근 아님)",
                          "특정심혈관질환(기타부정맥) = I49",
                          "심근병증(허혈성제외) = I42·I43 / 주요심장염증질환 = I30~41",
                          "허혈성심질환진단비 = I20·24·25 단독"]),
         ("롯데손해보험 5가지", ["특정심장질환Ⅰ = 급성심근경색 I21~23 · I46.0",
                          "특정심장질환Ⅱ = 협심증·허혈 + 주요심장염증 I30~41",
                          "특정 15대 = 판막 · 심근병증 · 빈맥 · 심부전",
                          "기타 심장부정맥 = I49 / 방실차단·전도장애 = I44·I45"]),
         ("삼성화재 · 메리츠화재", ["허혈성심장질환 6가지 = 급성기 I21 · 후속 I22 · 합병증 I23 · 협심증 I20 · 기타급성 I24 · 만성 I25",
                            "[메리츠] 심장질환진단Ⅰ → 허혈성 / Ⅱ → 급성심근경색 (별도 상품 병존)"])],
    ]
    for si, group in enumerate(SETS):
        s = prs.slides.add_slide(blank); header(s, name, 8 + si, "심혈관 담보 분류")
        box(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.5),
            '★ "특정Ⅰ·Ⅱ"는 회사마다 뜻이 다르다 — 라벨 말고 질병코드로 확인. '
            '빈맥(I47·48)과 부정맥(I49)은 별개.', 9.5, True, RED)
        for i, (comp, lines) in enumerate(group):
            x = Inches(0.5 + (i % 2) * 6.4); yy = Inches(1.95 + (i // 2) * 2.45)
            rect(s, x, yy, Inches(6.05), Inches(2.25), WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
            rect(s, x, yy, Inches(6.05), Inches(0.36), NAVY)
            box(s, x, yy, Inches(6.05), Inches(0.36), comp, 10.5, True, WHITE,
                PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            box(s, x + Inches(0.12), yy + Inches(0.42), Inches(5.8), Inches(1.75),
                "\n".join("· " + l for l in lines), 8.5, False, BLACK)
        footer(s, name, 8 + si, TOT - 1)

    prs.save(out)
    return out, state


if __name__ == "__main__":
    src, dst = sys.argv[1], sys.argv[2]
    out, st = build(src, dst)
    print(f"OK  {out}  (CI state = {st})")
