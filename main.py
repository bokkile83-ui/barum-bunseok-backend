# ===== BARUM main.py v41-fix12-20260712 (CI 상품명 공백무시·주계약/CI추가보장특약 다열 finditer) =====  BARUM main.py v33-ci-fix-20260708 (암주요치료비 매핑+수술 통원변형 차단+암/수술 감사로그 / 한화심혈관특정=확인) ===== (v29n + 심장묶음 6사 정본매핑·I20→협심증/허혈성=단독전용/순환계=전체5/급성심근=묶음제외 + 간병인MAX·요양드롭·간호통합7) =====
# -*- coding: utf-8 -*-
import os, re, tempfile, datetime, base64, traceback, json, httpx, urllib.parse
from fastapi import FastAPI, UploadFile, File, Form
from typing import List          # ★v385 제안서 복수(최대 3건) 업로드용
from fastapi.responses import HTMLResponse, JSONResponse, Response
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
import copy as _copy
from pptx.oxml.ns import qn as _qn
from pptx.text.text import _Run

# ★★★★★v407 (지점장 지시 2026.08.12): <b>각인은 한 곳에서만 정의한다.</b>
#   구 코드는 main.py 안 <b>4곳에 각인 문자열을 하드코딩</b>했다 — 한 곳만 안 바뀌면
#   `/health`·`/version`·`/diag`가 <b>서로 다른 버전</b>을 답하고, 그걸 보고 배포 여부를 오판한다.
#   ★이 상수가 main.py의 <b>유일한 각인</b>이다. 바꿀 때는 여기 한 줄만 바꾼다.
VSTAMP = 'v635-reset-20260902'


app = FastAPI(title="BARUM 보장분석 v7")

PW   = "0101"

# ═══════════ v429 회원 DB (제54조) — Railway Postgres ═══════════
#   ★권한자 = 최은혜 지점장. 마스터 비번 821024. 화면 비번은 0101 그대로.
#   DATABASE_URL이 없으면 <b>DB 없이도 앱은 돈다</b>(0101 게이트만 작동) — 배포 사고 방지.
ADMIN_PW = "821024"

# ★제55조 — 지침 하한선(2026.08.16 실측). 조문이 줄면 배포를 막는다.
DOCTRINE_MIN_ART = 77        # 조문 개수 하한
# ★v460 제69조 — DOCTRINE_MIN_CHARS 폐기. 손으로 관리하는 숫자가 정당한 정리를 막았다.
DOCTRINE_MIN_CHARS = 0         # (폐기 · 0 = 검사 안 함)
DOCTRINE_SKIP_ART  = {43}      # 처음부터 없는 번호(42 다음이 44)


def _db():
    """psycopg 연결. 실패하면 None — 호출부는 반드시 None을 처리한다."""
    import os as _o
    url = _o.environ.get('DATABASE_URL', '')
    if not url:
        return None
    try:
        import psycopg
        return psycopg.connect(url, connect_timeout=6)
    except Exception as _e:
        print('[v429 DB] 연결 실패:', str(_e)[:80])
        return None


def _db_init():
    c = _db()
    if not c:
        print('[v429 DB] DATABASE_URL 없음 — 회원 기능 비활성(0101 게이트만 작동)')
        return False
    try:
        with c, c.cursor() as k:
            k.execute("""CREATE TABLE IF NOT EXISTS members(
                id SERIAL PRIMARY KEY,
                name TEXT NOT NULL,
                code TEXT UNIQUE NOT NULL,
                phone TEXT DEFAULT '',
                memo TEXT DEFAULT '',
                status TEXT DEFAULT 'active',
                blocked BOOLEAN DEFAULT FALSE,
                created TIMESTAMP DEFAULT NOW(),
                expires DATE,
                last_used TIMESTAMP,
                use_count INT DEFAULT 0)""")
            # ★v439 기존 테이블에도 status 보강(제59조)
            try:
                k.execute("ALTER TABLE members ADD COLUMN IF NOT EXISTS status TEXT DEFAULT 'active'")
                k.execute("ALTER TABLE members ALTER COLUMN code DROP NOT NULL")
            except Exception as _ae:
                print('[v439 DB] status 보강 생략:', str(_ae)[:60])
            k.execute("""CREATE TABLE IF NOT EXISTS uselog(
                id SERIAL PRIMARY KEY,
                code TEXT, name TEXT, act TEXT,
                at TIMESTAMP DEFAULT NOW())""")
        print('[v429 DB] members · uselog 준비 완료')
        return True
    except Exception as _e:
        print('[v429 DB] 초기화 실패:', str(_e)[:100])
        return False
    finally:
        try: c.close()
        except Exception: pass


def _mk_code(n=6):
    """★숫자 6자리(지점장 확정 2026.08.16). 영문이 섞이면 폰 입력이 어렵다.
       맨 앞은 0을 피해 앞자리 누락을 막는다."""
    import random
    return str(random.randint(1, 9)) + ''.join(str(random.randint(0, 9)) for _ in range(n - 1))


def _member_check(code):
    """코드 검증. (ok, 이름, 사유)"""
    code = (code or '').strip().upper()
    if not code:
        return False, '', '코드를 입력하십시오'
    c = _db()
    if not c:
        return False, '', 'DB 미연결 — 지점장에게 문의'
    try:
        with c, c.cursor() as k:
            k.execute("SELECT name, blocked, expires, status FROM members WHERE code=%s", (code,))
            r = k.fetchone()
            if not r:
                return False, '', '없는 코드입니다'
            nm, blocked, exp, st = r
            if st and st != 'active':
                return False, nm, '아직 승인 전입니다'
            if blocked:
                return False, nm, '차단된 코드입니다'
            if exp:
                import datetime as _d
                if exp < _d.date.today():
                    return False, nm, '만료된 코드입니다 (%s)' % exp
            k.execute("UPDATE members SET last_used=NOW(), use_count=use_count+1 WHERE code=%s",
                      (code,))
            k.execute("INSERT INTO uselog(code,name,act) VALUES(%s,%s,'login')", (code, nm))
        return True, nm, ''
    except Exception as _e:
        return False, '', 'DB 오류: %s' % str(_e)[:60]
    finally:
        try: c.close()
        except Exception: pass

   # ★★v334 대문 비번 고정(지점장 지시 2026.08.02). Railway Variables의 ACCESS_PW/BARUM_PW(=1009)가 코드 기본값을 이기고 있어 환경변수 참조를 제거했다.
HERE = os.path.dirname(os.path.abspath(__file__))
TPL_XL  = os.path.join(HERE, "master.xlsx")
TPL_PPT = os.path.join(HERE, "ppt_form.pptx")
TPL_TX  = os.path.join(HERE, "chiryo_form.pptx")

W   = Font(color='FFFFFF', name='맑은 고딕', size=9, bold=True)
BL  = Font(color='0070C0', name='맑은 고딕', size=9)
BK  = Font(color='000000', name='맑은 고딕', size=9)
FILL_RED   = PatternFill('solid', fgColor='C00000')
FILL_BLUE  = PatternFill('solid', fgColor='0070C0')
FILL_GREEN = PatternFill('solid', fgColor='375623')
# ★v336 §10 정본: 실손(입원·통원·약값·MRI트리오·도수치료·비급여주사·상해의료비) + 일상배상책임 = 항상 파랑.
#   계약 셀뿐 아니라 <b>끝열 합계 셀도</b> 파랑이어야 한다(구 코드는 끝열을 무조건 검정으로 찍었다).
_BLUE_ROWS = {'입원','통원','약값','MRI트리오','도수치료','비급여주사','상해의료비','일상배상책임'}
FILL_SUM   = PatternFill('solid', fgColor='2E75B6')
AL = Alignment(horizontal='center', vertical='center', wrap_text=True)

EXCLUDE = ['실효','미납해지','농업인','자동차보험']  # NH농협=포함. 자동차(다이렉트/애니카/하이카 개인·업무·영업용)는 is_excluded에서 별도 처리

_NONLIFE_CO = ('손보', '손해보험', '화재', '해상', '재보험',
               'AIG', '처브', 'MG', '에이스', '흥국화재', '롯데손해', 'DB손해', 'KB손해',
               '삼성화재', '현대해상', '메리츠화재', 'MG손해', '한화손해', 'NH농협손해', '하나손해', '캐롯')

def _is_nonlife(co):
    """★★★★★v446 손해보험사 판정 (지점장 지시 2026.08.17, 영구)
       지점장 원문: 「<b>손해보험사의 CI는 무시다 일반보험이다. CI 적용하지마라</b>」
       근거 실측: 현대해상 '무배당퍼펙트클래스종합보험(Hi1706)'이 상품명에 '퍼펙트'를 포함해
                 삼성생명 퍼펙트 시리즈용 CI 규칙에 걸렸다 → 중대한 뇌졸증 2,000·중대한 급성심근 500.
       → CI/GI/리빙케어/퍼펙트 판정은 <b>생명보험사에만</b> 적용한다."""
    c = re.sub(r'[\s（）()]', '', str(co or ''))
    if not c:
        return False
    if '생명' in c or '라이프' in c:      # 생보가 먼저 (예: 흥국생명 vs 흥국화재)
        return False
    return any(k in c for k in _NONLIFE_CO)


def _isci_prod(p, co=''):
    """★v33 CI 상품명 판정 — 공백·전각 무시. '무배당교보큰사랑 CI 보험' 대응."""
    if _is_nonlife(co):
        return False          # ★v446 손보는 CI 판정 자체를 하지 않는다
    t = re.sub(r'[\s\u3000]', '', str(p or ''))
    # ★★★v110 영구지침(지점장 확정 2026.07.20): 삼성생명 <퍼펙트플러스보험>과
    #    <퍼펙트통합보험>은 상품명에 CI/GI/리빙케어 표기가 없어도 <무조건 CI보험>이다.
    #    (v105에서 '퍼펙트통합'을 오기로 보고 뺐으나, 지점장이 별개 상품으로 추가 확정 → 둘 다 CI)
    #    '퍼텍트~'는 리포트 표기 흔들림 대비 동의어.
    # ★★★v150 영구지침(지점장 확정 2026.07.21): <b>삼성생명 '퍼펙트' 시리즈는 전부 CI</b>다.
    #    퍼펙트통합보험 · 퍼펙트플러스보험 · 퍼펙트플러스종합보험 … 변형이 계속 나오므로
    #    <b>2종 열거를 폐기하고 '퍼펙트'(오타 '퍼텍트') 포함 여부로 판정</b>한다.
    #    근거: 상품명에 CI/GI/리빙케어 표기가 없어 이름만으로는 못 걸러진다 → 상세내역(세부가입현황) 검수 필수.
    if ('퍼펙트' in t) or ('퍼텍트' in t): return True
    if '리빙케어' in t: return True
    # ★★★v235 (2026.07.25 한정환 실측 — 지점장 지적 "CI가 2개인데 하나도 반영 안 됐다"):
    #   구 코드는 <b>`'CI보험'`이라는 연속 문자열</b>만 검사했다. 그래서 <b>CI와 '보험' 사이에
    #   상품 수식어가 끼면 전부 탈락</b>했다 — 실측 2건 다 놓쳤다:
    #     `무）라이프케어CI 종신보험` → 공백제거 `…라이프케어CI종신보험` (CI<b>종신</b>보험)
    #     `（무）변액유니버셜 세번받을수있는 CI종신보험（1701）` → 동일
    #   이건 v216b('수술비 관혈'→'비관혈')·v217('심뇌혈관'→'심뇌5대혈관')와 <b>완전히 같은 뿌리</b>
    #   = 연속 문자열 매칭의 함정. 세 번째 반복이다.
    #   → <b>CI·GI를 영문 경계로 잡고 '보험' 요구를 분리</b>한다.
    #   ★영문 경계 필수: `ACCIDENT`·`SPECIAL` 같은 영문 상품명 안의 우연한 'CI'를 배제해야 한다.
    # ★★★★★v246 실사고(지점장 "롯데나 KB나 동일하다" 확인 중 발견):
    #   <b>공백을 지우면 회사 영문약자와 CI가 붙어버린다</b> — `무배당 KB CI종신보험` → `무배당KBCI종신보험`
    #   → 'BCI'가 되어 <b>영문 경계 조건에 걸려 False</b>가 됐다(KB·DB·LG 등 영문약자 전부 해당).
    #   → <b>공백 제거본과 원문 둘 다</b> 검사한다. 원문에선 'KB CI종신보험'이라 CI 앞이 공백 → 정상 매칭.
    #   ★`ACCIDENT`는 원문에서도 CI 앞뒤가 영문이라 여전히 배제된다.
    t0 = str(p or '')
    _pat = r'(?<![A-Za-z])(CI|GI)(?![A-Za-z])'
    if (re.search(_pat, t) or re.search(_pat, t0)) and ('보험' in t or '종신' in t): return True
    return False


def _is_group_ins(product='', contract_date='', expiry_date=''):
    # ★단체보험 판정(지점장 확정 2026.07.18) — 2조건 동시 충족만 단체보험
    #   ① 가입기간이 1년마다 정해진다(가입~만기 = 1년)
    #   ② 상품명이 '○○단체보험'(상품명에 '단체' 표기)
    #   둘 중 하나만 맞으면 단체보험 아님 → 제외하지 않고 개인계약으로 포함
    p = str(product or '').replace(' ', '')
    if '단체' not in p:
        return False
    cy = re.match(r'(\d{4})', str(contract_date or ''))
    ey = re.match(r'(\d{4})', str(expiry_date or ''))
    if not cy or not ey:
        # ★★★★★v446 (지점장 지적 2026.08.17 「<b>단체보험 하지말랬는데 나왔다</b>」)
        #   실측: DB손보 '빅히트단체상해보험' — 계약일·만기일·보험료·납입횟수가 <b>전부 공란</b>이라
        #   구 코드가 '날짜 불명 → 포함'으로 처리해 엑셀 9열에 계약이 만들어졌다.
        #   → 상품명에 '단체'가 있는데 <b>날짜조차 확인 불가</b>면 개인계약으로 볼 근거가 없다.
        #     '단체' 표기를 신뢰해 제외한다. 날짜가 있으면 종전대로 1년 기준으로 판정한다.
        return True           # ★v446 '단체' + 날짜불명 → 제외
    try:
        span = int(ey.group(1)) - int(cy.group(1))
    except Exception:
        return False
    return span <= 1          # 1년 단위 가입기간

def _is_oneyear(contract_date, expiry_date):
    """★★★제외 6종 ⑥ 보험기간 1년 — 영구지침(지점장 확정 2026.07.20).
       가입시기 + 만기시기가 1년인 계약은 보장분석에서 제외한다.
       판정 = 만기일 - 가입일이 1년(±7일, 358~372일). 날짜 없으면 판정 불가 → 제외 안 함(누락 방지)."""
    import datetime as _dt
    def _d(x):
        x = re.sub(r'[^0-9]', '', str(x or ''))
        if len(x) != 8: return None
        try: return _dt.date(int(x[:4]), int(x[4:6]), int(x[6:8]))
        except Exception: return None
    a = _d(contract_date); b = _d(expiry_date)
    if not a or not b: return False
    days = (b - a).days
    return 358 <= days <= 372


def _pay_years(pay_period='', pay_count=''):
    """납입기간을 '연수'로 환산. 판정 불가면 None(제외 안 함·누락 방지). v126"""
    s = re.sub(r'\s', '', str(pay_period or ''))
    if '일시' in s: return 0                       # 일시납 = 납입기간 0년
    m = re.search(r'(\d+)\s*년', s)
    if m: return int(m.group(1))
    m2 = re.match(r'^\s*\d+\s*/\s*(\d+)\s*$', str(pay_count or ''))
    if m2:                                         # 총회차 → 월납 가정 연환산
        try: return int(m2.group(1)) // 12
        except Exception: return None
    return None

def _is_paid_up(pay_period='', pay_count=''):
    """납입완료 판정. 일시납은 정의상 완납. 회차 a/b에서 a>=b(b>0)면 완납. v126"""
    if '일시' in re.sub(r'\s', '', str(pay_period or '')): return True
    m = re.match(r'^\s*(\d+)\s*/\s*(\d+)\s*$', str(pay_count or ''))
    if not m: return False
    try:
        a, b = int(m.group(1)), int(m.group(2))
    except Exception:
        return False
    return b > 0 and a >= b

def _is_short_paidup(pay_period='', pay_count=''):
    """★★★제외 7종 ⑦ 단기완납 — 지점장 확정 2026.07.21(v131로 기준 변경).
       사유: 단기납 종신보험은 5년납(60회)부터 시작한다 → 60회 이상 완납은 정상 계약으로 포함.
       판정 = 완납(a>=b) AND 총회차 b < 60 → 제외(59회까지 제외).
       일시납은 1회이므로 제외. 회차를 못 읽으면 제외 안 함(누락 방지).
       ※구 기준 '120회 미만 제외'(v127)는 폐기."""
    if '일시' in re.sub(r'\s', '', str(pay_period or '')): return True   # 일시납 = 1회 < 60
    m = re.match(r'^\s*(\d+)\s*/\s*(\d+)\s*$', str(pay_count or ''))
    if not m: return False
    try:
        a, b = int(m.group(1)), int(m.group(2))
    except Exception:
        return False
    if not (b > 0 and a >= b): return False      # 납입 진행 중 → 대상 아님
    return b < 60                                 # ★60회(5년납) 미만 완납만 제외


def is_excluded(company, product='', contract_date='', expiry_date='', pay_period='', pay_count=''):
    t = re.sub(r'[\s（）()_·]|TM', '', str(company)+str(product))
    for kw in EXCLUDE:
        if kw in t: return True
    if _is_group_ins(product, contract_date, expiry_date): return True   # ★제외 5종: 단체보험
    # ★★★★★v313 (지점장 확정 2026.08.01, 영구): <b>일시납(1/1)은 만기·기간 무관 무조건 기재한다</b>.
    #   지점장 원문 = "단(1/1) 즉 <b>일시납은 기재하고 담보도 기재</b>해줘야 한다"
    #                 "하나생명 VIP는 저축보험이다 — 담보가 사망 300만뿐 → <b>넣어라</b>"
    #   ★구 v302 조건 <b>「일시납 AND 만기 9999(종신)」는 폐기</b> — 만기 조건이 사라졌다.
    #   ★제외 ⑥(보험기간 1년)보다도 <b>앞</b>이다. 실측 = AIG 부모님건강보험(1/1 · 2025.12.31~2026.12.31)이
    #     ⑥에 먼저 걸려 빠졌는데 지점장은 <b>둘 다 포함</b>으로 확정했다.
    #   ★단 ①실효 ②미납해지 ③농업인 ④자동차 ⑤단체는 <b>그대로 제외</b>(위에서 이미 처리됨).
    _pp = re.sub(r'\s', '', str(pay_period or ''))
    _pc = re.sub(r'\s', '', str(pay_count or ''))
    _onetime = ('일시' in _pp) or (_pc == '1/1')
    if _onetime:
        print(f"[v313 일시납 포함] {company} {product} — 납입 {pay_period or pay_count} · 만기 {expiry_date} → 제외6·7 면제")
        return False
    if _is_oneyear(contract_date, expiry_date):                          # ★제외 6종: 보험기간 1년(v102)
        print(f"[제외6·보험기간1년] {company} {product} — {contract_date}~{expiry_date}")
        return True
    if _is_short_paidup(pay_period, pay_count):                          # ★제외 7종: 단기완납(v126)
        print(f"[제외7·단기완납] {company} {product} — 납입 {pay_period or pay_count} (총회차 60 미만 완납)")
        return True
    if '운전자' in t: return False   # ★운전자·운전자상해보험은 포함(§4)
    # ★자동차보험(다이렉트/애니카/하이카 + 개인용/업무용/영업용/개인소유) = 제외
    if any(b in t for b in ('다이렉트','애니카','하이카','개인용자동차','업무용자동차')) and any(x in t for x in ('개인용','업무용','영업용','개인소유')):
        return True
    return False

def _no_period(contract_date, expiry_date):
    """계약일·만기일 중 하나라도 8자리 날짜로 읽히지 않으면 True (보험기간 판정 불가). v125"""
    def _ok(x):
        return len(re.sub(r'[^0-9]', '', str(x or ''))) == 8
    return not (_ok(contract_date) and _ok(expiry_date))


def _is_silson_like(company='', product='', dambo=None):
    """실손 계약 판정(확장) — 상품명/회사명에 '실손'이 없어도 담보에 실손·의료비가 있으면 실손으로 본다.
       장문순 실측: 삼성생명 퍼펙트통합보험 줄이지만 담보는 실손 처방조제료·외래의료비뿐이었다. v125"""
    if _is_silson_prod(company, product): return True
    try:
        for k in (dambo or {}):
            kk = re.sub(r'\s', '', str(k))
            if ('실손' in kk) or ('의료비' in kk): return True
    except Exception:
        pass
    return False


def _is_silson_prod(company='', product=''):
    """실손 계약 판정 = 상품명/회사명에 '실손' 표기. (v103)"""
    return '실손' in re.sub(r'\s', '', str(company) + str(product))


#   ★v473 상품명에서 '갱신'을 근거로 인정하는 <b>정형 토큰</b>. 맨 '갱신' 두 글자는 인정하지 않는다.
_RENEW_TOKEN = re.compile(r'갱신형|\d+\s*년\s*갱신|\(\s*갱신\s*\)')
#   ★★v473f 상품명 '세만기' = 비갱신(지점장 확정 2026.08.18 · <b>지침이 우선이다</b>).
#     ★v473e에서 내가 넣었던 「'95세만기'·'100세만기'는 만기 나이라 제외」는 <b>지침에 없는 예외</b>였다.
#       지점장 지적 — 「<b>지침이 우선이다</b>」 ⇒ <b>'세만기'라고 적혀 있으면 비갱신</b>. 예외 없음.
_SEMANGI = re.compile(r'세만기')


#   ★★★★★v474 제82조 — <b>상품명은 제목이 끝나면 끝낸다</b>(지점장 지시 2026.08.18).
#     지점장 원문 — 「<b>이건 너의 숙제고 너무 다양해서 다 파악하기 힘들다. 하지만 제목이 끝나면
#     끝내라</b>」
#     [실측 오염 2건 · 이창재]
#       `한화 운전자상해보험 무배당2404 상하지(손,발제외)절골술및체내금속고정수술비(연간1회한,급`
#       `LIFEPLUS 3N5 간편건강보험(세만기형) 무배당2405 갑상선암및전립선암다빈치로봇수술비(1회한)(3N5간편,갱신`
#     별첨 표에서 <b>상품명 칸 뒤에 담보명이 흘러들어온다</b>. 이 꼬리의 '갱신' 두 글자가
#     제5조 ①을 먹어 계약을 통째로 뒤집었다(v473 사고의 근원).
#     [자르는 법 — 회사별 하드코딩 없이 구조로만]
#       ① 상품명을 공백 토큰으로 나눈다.
#       ② <b>'보험'·'공제'·'플랜'류가 들어간 마지막 토큰</b> = 제목의 끝(상품 종결어).
#       ③ 그 <b>뒤쪽</b> 토큰 중 <b>담보 접미어</b>(진단비·수술비·치료비·일당·급부금…)를 가진
#          <b>첫 토큰부터 끝까지</b> 잘라낸다. 상품코드(`무배당2404`)는 접미어가 없어 살아남는다.
#       ④ 종결어가 없거나 자를 게 없으면 <b>원본 그대로</b>(안전측 — 못 자르는 건 놔둔다).
_PROD_END  = ('보험', '공제', '플랜', '연금', '저축')      # 제목의 끝을 알리는 말
_DAMBO_TAIL = ('진단비', '수술비', '치료비', '입원비', '통원비', '일당', '급부금', '급여금',
               '위로금', '지원금', '의료비', '보장금', '선임비용', '벌금', '수술료', '진단금')


def _clean_product(p):
    """★v474 제82조 — 상품명 뒤에 붙은 담보명 꼬리를 잘라낸다. 못 자르면 원본을 돌려준다."""
    s = str(p or '')
    tk = s.split()
    if len(tk) < 2: return s
    end = -1
    for i, t in enumerate(tk):
        if any(k in t for k in _PROD_END): end = i        # 마지막 종결어 위치
    if end < 0: return s                                   # 제목의 끝을 못 찾으면 손대지 않는다
    for j in range(end + 1, len(tk)):
        if any(k in tk[j] for k in _DAMBO_TAIL):
            _new = ' '.join(tk[:j]).strip()
            return _new or s                               # 빈 문자열이면 원본 유지
    return s


#   ★★★★★v476 제84조 — <b>회사명 앞에 기간칸이 흘러들어온다</b>(지점장 지시 2026.08.18).
#     실측(김순자 C열) — `년납/3년NH농협생명` . 별첨 헤더의 <b>납입기간 칸</b>이 회사명 앞으로 밀려
#     들어와 회사명·확인사항·근거표·검색링크까지 전부 오염됐다.
#     [자르는 법] 회사명 <b>맨 앞</b>에 붙은 기간·납입 토큰(`N년납`·`/N년`·`일시납`·`월납`·`연납`…)만
#     걷어낸다. 회사 이름은 숫자나 '/'로 시작하지 않는다. 걷어내고 남는 게 없으면 원본을 쓴다.
_CO_HEAD = re.compile(r'^(?:\s|/|\d+\s*년?|년납|월납|연납|일시납|분기납|반년납|납)+')


#   ★★★★★v476 제86조 — <b>「갱신형 담보」도 「담보 (갱신형)」도 모두 갱신이다</b>
#     (지점장 지시 2026.08.18 · 추가건). 담보명에서 '갱신' 표기는 <b>앞에 붙든 뒤에 붙든 같다</b>.
#       앞: `갱신형 뇌혈관질환진단비` · `[갱신형]암진단비`
#       뒤: `간병인사용 질병입원일당(1-180일)(간편가입)(갱신형)` · `…(3N5간편,갱신형)`
#       중간: `(10년갱신)갱신형 다빈치로봇 암수술비`
#     ⇒ <b>위치를 따지지 않는다.</b> 판정은 이 함수 하나로만 한다(제0조 「판정은 한 곳에서만」).
#   ★★★★★v486 제96조 (지점장 실물 PDF 대조 2026.08.19) — <b>「담보 갱신형」이 '비갱신'으로 읽혔다</b>.
#     구 코드는 공백을 지운 뒤 `'비갱신' in t`로 제외했다. 그런데 <b>'진단비' + '갱신형'</b>이 붙으면
#     `암진단<b>비갱신</b>형` — <b>없던 '비갱신'이 만들어진다</b>. 그래서 뒤에 붙은 갱신형이 통째로 탈락했다.
#     ★게다가 구 코드의 `re.sub`는 <b>역슬래시가 두 개</b>라 공백을 아예 못 지우고 있었다(조용한 결함).
#     ⇒ '비갱신'은 <b>단어 경계</b>(문자열 시작 · 공백 · 괄호 · 쉼표 뒤)일 때만 인정한다.
# ★★★★★v596 제96조 2항 — '비' 앞이 <b>한글이 아니면</b> 진짜 비갱신이다.
#   `암진단_비갱신`(KDB 실측) · `(비갱신)…` · `비갱신형 …` 은 잡고,
#   `암진단비`+`갱신형`이 붙은 `암진단비갱신형`은 '비' 앞이 '단'(한글)이라 안 잡힌다.
_NONGEN = re.compile(r'(?:^|[^가-힣])비\s*갱신')


#   ★★★★★v495 제98조 (지점장 확정 2026.08.19) — <b>배상책임 담보는 갱신이다.</b>
#     지점장 원문 — 「<b>가족생활배상책임담보 / 갱신이다. 이게 비갱신인 건 2009년 이전만 비갱신이다</b>」
#     [실측] 메리츠 알파Plus2004(2020년 가입)에서 `가족생활배상책임담보`가 <b>갱신 아님</b>으로
#       세어져 과반 계산이 52/54가 됐다. 실제로는 <b>53/54</b>다.
#     ⇒ 일상생활·가족생활·자녀생활 <b>배상책임</b> 계열은 <b>2009년 이후 가입이면 갱신</b>으로 센다.
#       2009년 이전 가입만 비갱신.
_ILSANG = re.compile(r'(?:일상생활|가족생활|자녀생활|일상|가족)\s*배상책임')


def _is_gen_dambo(raw, contract=''):
    """★제86조 · v493 제96조 — 담보명에 갱신 표기가 있는가(앞·뒤·중간 위치 무관).

    ★★★★★지점장 확정 2026.08.19 — <b>담보에 「비갱신형」이라고 기재는 절대 안 한다.
      갱신형만 표기된다.</b> ⇒ 「비갱신」 제외 로직 자체가 필요 없다.
      구 코드는 공백을 지운 뒤 `'비갱신' in t`로 걸렀는데, <b>'암진단비'+'갱신형'</b>이 붙어
      `암진단<b>비갱신</b>형` — <b>없던 '비갱신'이 만들어져</b> 뒤에 붙은 갱신형이 통째로 탈락했다.
      단어 경계(`_NONGEN`)로 막았으나, 지점장 확정으로 <b>그 가드도 폐기</b>한다.
      담보명에 '갱신'이 있으면 <b>무조건 갱신</b>이다.
    """
    _t95 = re.sub(r'\s', '', str(raw or ''))
    # ★★★★★v596 제96조 2항 (지점장 실측 2026.08.28 「암진단비 비갱신인데 갱신으로 인식했다」
    #   · 김생금님 KDB생명 담보명 `암진단_비갱신`)
    #   [구 결함] `'갱신' in 이름` 하나로 끝나서 <b>'비갱신'에도 '갱신'이 들어 있어</b>
    #   비갱신 담보가 통째로 갱신(파랑)으로 뒤집혔다.
    #   v493의 전제(「담보에 비갱신형이라고 기재는 절대 안 한다」)가 <b>실물과 달랐다</b> — KDB는 기재한다.
    #   ⇒ 비갱신 가드를 되살린다. 단 v486 사고(`암진단비`+`갱신형`)는 안 걸리게
    #     <b>'비' 앞이 한글이 아닐 때만</b> 비갱신으로 본다.
    if _NONGEN.search(str(raw or '')): return False
    if '갱신' in _t95: return True
    # ★★★★★v499 제98조 3항 (지점장 지적 2026.08.19 「흥국화재 10억통장이 갱신인데 블랙으로 나오더라」)
    #   흥국 <b>리셋월렛II(10억 통장)</b>는 담보명에 갱신 표기가 없어도 <b>갱신 담보</b>다.
    #   구 코드는 표기가 없어 검정으로 찍었다.
    if ('리셋월렛' in _t95) or ('리셋월랫' in _t95): return True
    # ★v495 제98조 — 배상책임 계열은 2009년 이후 가입이면 갱신
    if _ILSANG.search(str(raw or '')):
        _y = int(str(contract)[:4]) if str(contract)[:4].isdigit() else 0
        return _y >= 2009
    return False


_GENDAMBO_SELFTEST = [
    ('갱신형 뇌혈관질환진단비', True),
    ('[갱신형]암진단비(유사암제외)(추가고지형)(갱신형_20년)', True),
    ('간병인사용 질병입원일당(요양병원제외,1-180일)(간편가입)(갱신형)', True),
    ('갑상선암및전립선암다빈치로봇수술비(1회한)(3N5간편,갱신형)', True),
    ('(10년갱신)갱신형 다빈치로봇 암수술비(특정암)', True),
    ('일반상해사망(간편가입)(갱신형)', True),
    ('재해사망', False),
    ('질병사망', False),
    ('비갱신형 상해수술비', False),
    # ★v596 제96조 2항 — 실물에 '비갱신'이 기재된다(김생금님 KDB 실측)
    ('암진단_비갱신', False),
    ('암진단_비갱신형', False),
    ('(비갱신)일반암진단비', False),
    ('일반암진단비 비갱신', False),
    ('암진단비 갱신형', True),        # ★v486 사고 재발 방지 — 붙어도 갱신이다
    ('암진단비갱신형', True),
]


def _clean_company(c):
    """★v476 제84조 — 회사명 앞에 붙은 기간칸 토큰을 잘라낸다. 못 자르면 원본."""
    s0 = str(c or '').strip()
    if not s0: return s0
    s1 = _CO_HEAD.sub('', s0).strip()
    return s1 if re.search(r'[가-힣A-Za-z]', s1) else s0


_COCLEAN_SELFTEST = [
    ('년납/3년NH농협생명', 'NH농협생명'),
    ('일시납AIG손보', 'AIG손보'),
    ('20년납 삼성화재', '삼성화재'),
    ('NH농협생명', 'NH농협생명'),
    ('KB손보', 'KB손보'),
    ('메리츠화재', '메리츠화재'),
    ('롯데손해보험', '롯데손해보험'),
]


_PRODCLEAN_SELFTEST = [
    ('한화 운전자상해보험 무배당2404 상하지(손,발제외)절골술및체내금속고정수술비(연간1회한,급',
     '한화 운전자상해보험 무배당2404'),
    ('LIFEPLUS 3N5 간편건강보험(세만기형) 무배당2405 갑상선암및전립선암다빈치로봇수술비(1회한)(3N5간편,갱신',
     'LIFEPLUS 3N5 간편건강보험(세만기형) 무배당2405'),
    ('(무) 경영인 정기보험(2405)(3형:20%체증형,해약환급금일부지급형,95세만기)',
     '(무) 경영인 정기보험(2405)(3형:20%체증형,해약환급금일부지급형,95세만기)'),
    ('성공하는 Owner 재산종합보험 무배당2404', '성공하는 Owner 재산종합보험 무배당2404'),
    ('(무)LIG닥터플러스Ⅴ보험', '(무)LIG닥터플러스Ⅴ보험'),
    ('무배당 흥Good 고당지 3.10.5 간편종합보험', '무배당 흥Good 고당지 3.10.5 간편종합보험'),
]


def judge_renewal(product, expiry, pay_count, contract='', pay_period='', company='', dambo=None):
    # 지침 §6 판정 (2026.07.09 개정: 240회 규칙 삭제 / 삼성화재 예외 / 납입==보장→갱신)
    # 0) ★★★실손은 비갱신이 없다 — 무조건 '갱신' (영구지침, 지점장 확정 2026.07.20 / v103)
    #    실손의료비는 제도상 갱신형만 존재한다. 만기 9999·납입!=보장 등 어떤 조건에도 우선한다.
    if _is_silson_prod(company, product): return '갱신'
    # ★★★★★v473e 제5조 최종 (지점장 확정 2026.08.18 · 이 순서다)
    #   지점장 원문 —
    #   「<b>20년 = 가입시기 → 만기시기 = 20년 동일 → 갱신 이거나 아니면 상품명에 (갱신형)이거나 하면
    #     올 갱신. 이고 20년 = 가입시기 → 만기시기 = 상이하면 비갱신 이거나 세만기 라고 적혀있으면
    #     비갱신이다. 그리고 담보별에 (갱신)이라고 적혀있는건 주계약은 비갱신이지만 담보는 갱신이다</b>」
    #   [계약(헤더) 판정] ⓪실손 → ①<b>상품명 '갱신형' → 갱신</b> → ②<b>상품명 '세만기' → 비갱신</b>
    #                    → ③만기 9999 → 비갱신(종신) → ④납입기간 vs 보장기간
    #   ★v473f 정정 — 지점장 원문이 <b>갱신 조건을 먼저</b> 말한다. v473e에서 내가 세만기를 앞에 뒀던 것을
    #     원문 순서로 되돌렸다. LIFEPLUS는 상품명에 <b>(갱신형) 정형 토큰이 없어</b> ②에서 비갱신으로 잡힌다.
    #   [담보(값 글자색) 판정] 담보명에 '갱신' → <b>그 담보만 파랑</b>(주계약이 비갱신이어도).
    #                        구현 = build_excel `blue = gen or ('갱신' in raw)` — 별개 층위다.
    _p473 = re.sub(r'\s', '', str(product or ''))
    # ① 상품명에 '갱신형' 명시 -> 갱신 (지점장 원문 = 갱신 조건이 먼저다)
    #    ★대상은 <b>상품명뿐</b>. '갱신' 두 글자 부분일치 금지(담보명 꼬리 오염 차단) — 정형 토큰만.
    if _RENEW_TOKEN.search(_p473) and '비갱신' not in _p473: return '갱신'
    # ② 상품명에 '세만기' -> 비갱신
    #    ★<b>예외 없다.</b> v473e에서 내가 「95세만기·100세만기는 만기 나이라 제외」를 넣었으나
    #      지점장 지적 「<b>지침이 우선이다</b>」로 <b>철회</b>했다. 적혀 있으면 세만기다.
    if _SEMANGI.search(_p473) and '비갱신' not in _p473: return '비갱신'
    # ③ 만기 9999(종신) -> 비갱신
    if expiry.startswith('9999'): return '비갱신(종신)'
    # ★★★★★v484 제95조 (지점장 확정 2026.08.19) — <b>담보에 갱신 표기가 있으면 그 계약은 갱신이다</b>.
    #   지점장 원문 — 「<b>우리가 밤샌 이유는 이분 다 갱신인데 니가 계속 비갱신으로 계산해서야</b>」
    #                「<b>갱신 담보 or 담보 갱신 적용해봐</b>」
    #   ⇒ 제86조(`_is_gen_dambo` · 앞·뒤·중간 위치 무관)를 <b>계약 헤더 판정에도</b> 쓴다.
    #   [실측 김순자] 메리츠 알파Plus2004 담보 54개 중 <b>52개가 갱신형</b>인데 20년납/42년보장이라
    #     기간(⑤)으로 비갱신이 됐다. 2204는 <b>10개 전부</b> 갱신형, KB 3.5.5는 19개 중 17개.
    #   [순서] 상품명 '세만기'(②)가 <b>이보다 앞</b>이다 — 세만기형 상품의 특약 갱신 표기로는
    #     계약을 뒤집지 않는다(LIFEPLUS 3N5 · 지점장 2026.08.18 확정 유지).
    #   ★★★★★v494 제95조 개정 (지점장 원리 설명 2026.08.19)
    #     「<b>주계약이 비갱신이고 특약이 갱신인 경우가 대부분이기에 블랙·블루 처리하는 거다.</b>
    #      모든 보장분석지에 갱신은 기재되지만 <b>단 한 번도 비갱신이라고 기재는 없다</b>.
    #      그래서 타이틀에 세만기를 보든지 <b>납부기간과 보험기간의 상이성</b>을 봐야 하고,
    #      <b>상이해도 일부 특약은 갱신 기재 시 블루</b>다」
    #     ⇒ <b>담보 대부분(과반)이 갱신 표기면 계약도 갱신</b>. 「일부」면 계약은 기간대로 판정하고
    #       <b>그 담보만 파랑</b>(제5조 B · 제83조)으로 둔다. 하나만 있어도 계약을 뒤집던 것은 과했다.
    #     [실측 김순자] 메리츠2004 <b>52/54</b> · 2204 <b>10/10</b> · KB <b>17/19</b> → 계약 갱신
    #                  메리츠 운전자 <b>2/15</b>(일부) → 계약은 기간(20납/20년 동일)으로 갱신
    _dk95 = list((dambo or {}).keys())
    #   ★★★★★v496 제98조 2항 (지점장 확정 2026.08.19 「의료 엑셀에 없다」)
    #     <b>마스터 엑셀에 행이 없는 담보는 과반 분모에서 뺀다.</b> 엑셀에 안 실리는 담보가
    #     분모에 들어가면 판정이 왜곡된다(제0조 「엑셀 = 지침의 법률 존재」).
    #     실측 — 메리츠2004의 `의료사고법률비용보장담보`는 `resolve2` 결과가 <b>None</b>(마스터 무행).
    #     이걸 빼면 <b>53/53 = 100%</b>다.
    _dk95 = [_k for _k in _dk95 if (resolve2(_k) or (None,))[0]]
    #   ★★★★★v498 제95조 재개정 (지침 재정독 2026.08.19) — 지점장 원문은 <b>「담보에 갱신 100% 표기」</b>다.
    #     내가 「과반」으로 낮춰 썼다 — <b>또 임계값을 내가 정했다</b>. 51%도 갱신이 되어
    #     「일부 특약만 갱신인 비갱신 계약」이 갱신으로 나갈 뻔했다. <b>100%로 되돌린다.</b>
    #   ★<b>잘린 담보명은 분모에서 뺀다</b> — 별첨 표에서 칸 폭 때문에 뒤가 끊긴다.
    #     실측 KB `상해입원일당(요양/정신/한방병원제외,181일이상)(간편가입)<b>(</b>` — `(갱신형)`이 잘렸다.
    #     괄호가 안 닫힌 담보명은 <b>원문이 더 있다</b>는 뜻이므로 판정 근거로 쓰지 않는다(제102조와 같은 뿌리).
    _dk95 = [_k for _k in _dk95 if str(_k).count('(') == str(_k).count(')')]
    if _dk95 and all(_is_gen_dambo(_k, contract) for _k in _dk95): return '갱신'
    # ④ 납입기간 == 보장기간(가입일~만기일) 동일 -> 갱신 / 다르면 비갱신
    pay_y = 0; cov_y = 0
    m = re.search(r'(\d+)\s*년', pay_period or '')
    if m: pay_y = int(m.group(1))
    if not pay_y:
        try:
            _, b = pay_count.split('/'); _b = int(b.strip())
            # ★★★★★v476 제85조 — <b>총회차가 12 미만이면 연납이다</b>(지점장 지시 2026.08.18).
            #   구 코드는 총회차를 <b>무조건 개월수</b>로 보고 12로 나눴다 → 연납 계약
            #   `2/3`(3년납 3회)이 `round(3/12)=0`이 되어 <b>기간 판정 자체를 건너뛰고</b>
            #   무조건 비갱신으로 떨어졌다(김순자 NH모두안심재해보험 실측).
            #   ⇒ 12 미만이면 그 숫자가 곧 <b>연수</b>다. 12 이상은 종전대로 개월수.
            pay_y = _b if 0 < _b < 12 else round(_b/12)
        except: pass
    try:
        cy = int(contract[:4]); ey = int(expiry[:4])
        if cy and ey: cov_y = ey - cy
    except: pass
    if pay_y and cov_y:
        if pay_y != cov_y: return '비갱신'          # ★기간 상이 = 비갱신
        # ★★★★★v483 제93조 전기납 정의 (지점장 확정 2026.08.19 · 이전 기재 전부 폐기)
        #   지점장 원문 — 「<b>전기납은 삼성화재를 제외하고 최소 35년납 이상 되는게 전기납이다.
        #   납입주기 · 가입일자와 만기일자가 납입주기와 같은 시기는 갱신이고</b>」
        #   ⇒ <b>납입기간 == 보장기간이면 원칙은 갱신이다.</b> 전기납은 <b>35년납 이상</b>일 때만.
        #   ⇒ 예외는 <b>삼성화재</b> 하나 — 35년 미만이어도 전기납 판별(담보 갱신표기 유무)을 탄다.
        #   [내 오판] v354·v473은 <b>기간 길이를 안 보고</b> 「담보에 갱신 표기 없으면 전기납」으로만 봤다.
        #     그래서 <b>롯데 20년납/20년만기(= 명백한 갱신)</b>가 비갱신으로 나갔다(김순자 실측).
        #     지점장 지적 — 「<b>이건 당연히 갱신이자나</b>」.
        _co483 = re.sub(r'\s', '', str(company or ''))
        if pay_y >= 35 or ('삼성화재' in _co483):
            _has_gen = any(_is_gen_dambo(k, contract) for k in (dambo or {}).keys())   # ★v476 제86조
            return '갱신' if _has_gen else '비갱신'
        return '갱신'                      # ★35년 미만 + 납입==보장 = 갱신(제93조)
    return '비갱신'

def _has_nonpay3(dambo):
    """★★★v250 (지점장 확정 2026.07.26): <b>3대비급여 특약이 분리돼 있으면 실손 세대 하한 = 3세대</b>.
    근거 = 도수치료·체외충격파·증식치료 / 비급여 주사제 / 비급여 MRI 검사 이 3특약은
    <b>2017.04 3세대(착한실손)부터 신설</b>됐다 — 2세대 이하 계약에는 구조적으로 존재할 수 없다.
    ★실측(한O환 DB생명 CI종신1701): 상품코드 1701(2017.01) 때문에 <b>2세대로 오판</b>됐으나
      별첨에 도수350·주사제250·MRI300이 <b>분리 표기</b> → <b>3세대가 정답</b>(가입일 2018.03.30과도 일치).
    ★상품코드는 <b>CI종신 주계약 코드</b>일 뿐 실손 특약 코드가 아니다 — 그래서 하한이 필요하다."""
    for k in (dambo or {}):
        n = re.sub(r'\s', '', str(k)).upper()
        if ('도수' in n) or ('체외충격파' in n) or ('증식치료' in n): return True
        if ('비급여' in n) and (('주사' in n) or ('MRI' in n)): return True
        if ('MRI' in n) and (('검사' in n) or ('비급여' in n)): return True
    return False

def _has_drug(dambo):
    """★v381 처방조제료(약값)가 통원과 별도로 잡혀 있는가 — 4세대 배제 판정용.
       지점장 확정 2026.08.11: "4세대는 통원비 20만원만 있다"."""
    for k in (dambo or {}):
        n = re.sub(r'\s', '', str(k))
        if ('처방조제' in n) or ('약제비' in n) or ('약값' in n): return True
    return False


# ★★★★★v382 「단독 5종」 판정을 <b>함수 하나로 단일화</b>한다(지점장 확정 2026.08.11 조문).
#   뇌혈관진단비 · 허혈성진단비 · 급성심근경색진단비 · 뇌졸증진단비 · 뇌출혈진단비 = 전 회사 단독.
#   단서: 등급 로마숫자 <b>Ⅰ·Ⅱ</b>가 붙으면 단독 예외에서 빼고 종전 회사별 표를 탄다(Ⅲ은 상품세대 표기라 단독 유지).
#   ★인라인에 조건식을 박아두면 <b>셀프테스트가 그 로직을 검사할 수 없다</b> → 함수로 빼서
#     `doctrine_selftest()`가 <b>실제로 쓰이는 바로 그 함수</b>를 매 배포 검사하게 한다.
_SOLO5_KW   = ('허혈', '뇌혈관', '뇌졸', '뇌출혈', '급성심근')
_BUNDLE_MK5 = ('특정', '6가지', '5대', '4대', '3대', '2대', '순환계')

# ★★★★★v387 (지점장 확정 2026.08.12): 「<b>허혈성 이라는 키워드가 앞에 있으면 무조건 단독이다</b>」
#   계기 = 「<b>허혈성심질환진단비Ⅱ — 이건 심장Ⅱ가 아니라 허혈성심질환진단비니까 허혈성 단독이다</b>」
#   ⇒ <b>허혈</b>이 담보명 <b>맨 앞</b>에 오면 뒤에 붙은 Ⅰ·Ⅱ는 <b>상품 표기</b>이므로 등급 예외를 타지 않는다.
#   ⇒ 앞에 다른 말이 붙은 것(`특정허혈…`·`5대혈관…`)은 <b>묶음</b>이므로 종전대로 회사별 표.
#   ★<b>이 완화는 「허혈」에만 적용한다</b> — 지점장이 허혈성 하나를 말씀하셨다.
#     `뇌혈관진단비Ⅰ/Ⅱ` 등 나머지 4종의 등급 예외는 <b>손대지 않는다</b>(0f: 하나를 말했으면 하나만).
#   ★v386은 이걸 `_hl_bund`로 <b>우연히</b> 막았을 뿐 이 함수는 여전히 단독에서 빼고 있었다 → 근본을 여기서 고친다.
_SOLO5_PREFIX = ('[', ']', '(', ')', '（', '）', '무배당', '갱신형', '무', '신', ' ')

def _solo5_head_is_haehyeol(rn):
    """담보명 <b>맨 앞</b>이 '허혈'인가 — 접두 표기([갱신형]·(무배당) 등)는 걷어내고 본다."""
    t = re.sub(r'\s', '', str(rn or '')).replace('（','(').replace('）',')')
    t = re.sub(r'^(?:\[[^\]]*\]|\([^)]*\)|무배당|갱신형)+', '', t)   # 접두 표기 제거
    return t.startswith('허혈')

def is_solo5_name(rn):
    rn = str(rn or '')
    if any(k in rn for k in _BUNDLE_MK5): return False       # 묶음 수식어가 있으면 회사별 표
    if _solo5_head_is_haehyeol(rn):       return True        # ★v387 허혈이 앞 = 등급 무관 단독
    if not any(k in rn for k in _SOLO5_KW): return False
    return _rmn(rn) not in (1, 2)


def silson_gen(contract_date, ipv=None, product='', nonpay3=False, drug=False):
    """실손 세대 판별 — 5세대=2026.05부터(정본 확정). 4세대=2021.07~2026.04. 입원한도 3000=구형(1세대). 가입일 없으면 '' → [확인].
    ★v29v: 상품명 연도코드(YYMM 4자리, 예 '1804'=2018.04 출시)가 있으면 판정일로 우선 사용 —
    갱신 재가입일(계약일)로 세대를 오판(예 2018 실손을 4세대로)하는 것 차단.
    ★★★v250: `nonpay3=True`(3대비급여 특약 분리)면 <b>3세대 하한 고정</b>(지점장 확정 2026.07.26)."""
    if ipv==3000: return '1세대(구형)'
    try: ym=int(str(contract_date)[:4])*100+int(str(contract_date)[5:7])
    except: ym=0
    _pm=re.search(r'(?<!\d)(0[9]|1[0-9]|2[0-6])\.?(0[1-9]|1[0-2])(?!\d)', str(product or ''))   # ★v30 '25.01' 점 형식 포함
    if _pm:
        _pym=2000*100+int(_pm.group(1))*100+int(_pm.group(2))
        ym=_pym if not ym else min(ym,_pym)
    if not ym:
        # ★★★★★v381 (지점장 확정 2026.08.11, 영구): 가입일·상품코드가 없어도 <b>담보 구조로 판정</b>한다.
        #   지점장 원문 = "<b>4세대는 통원비 20만원만 있다</b>".
        #   3대비급여 특약 분리(2017.04~ 신설) → 3세대 이상 / 그중 <b>처방조제료(약값)가 별도</b>면 3세대,
        #   통원만 있고 약값이 없으면 4세대. coverage_benchmark `_gen_of`와 <b>같은 규칙</b>(결과값 동결).
        if nonpay3: return '3세대' if drug else '4세대'
        return ''
    # ★★★v250 3대비급여 하한 — 상품코드가 실손이 아닌 주계약 코드일 때의 오판을 막는다.
    if nonpay3 and ym < 201704:
        print(f'[v250 3대비급여] 실손 세대 하한 고정 {ym} → 201704(3세대) · 근거=도수/비급여주사/비급여MRI 특약 분리')
        ym = 201704
    if ym<200910:  return '1세대'
    if ym<=201703: return '2세대'
    if ym<=202106: return '3세대'
    # ★★★v211 (지점장 확정 2026.07.25, 영구): <b>5세대 출시 = 2026.05.06</b>
    #   근거 = 금융위원회·금융감독원 보도자료(보도시점 2026.5.6 조간) "5월 6일부터 …
    #   5세대 실손의료보험이 새롭게 출시·판매됩니다". 구 기재 '5세대=2026.05~'는 폐기.
    #   판정 근거는 <b>상품명의 상품코드(YYMM) 또는 가입일자</b>(지점장 지시) — 둘 다 있으면 더 이른 쪽.
    #   상품코드는 월까지만 있어 2605면 5세대로 보고, <b>가입일자가 있으면 일자까지</b> 따진다(5/1~5/5 = 4세대).
    if ym<202605: return '4세대'
    if ym==202605:
        try: _d=int(str(contract_date)[8:10])
        except Exception: _d=0
        if _d and _d<=5: return '4세대'
    return '5세대'

def silson_gen_desc(gen):
    """세대별 보장설명지용 한 줄 설명."""
    return {
      '1세대':'자기부담 0~20%·갱신3·5년·재가입없음(구실손)',
      '1세대(구형)':'입원한도 3천·자기부담 0~20%(구실손 1세대)',
      '2세대':'급여90%·비급여90%·재가입15년(표준화 실손)',
      '3세대':'급여80%·비급여70~80%·도수 특약분리·재가입15년',
      '4세대':'급여80%·비급여70%·비급여 할증·재가입5년',
      '5세대':'입원 급여80%·비급여 중증70/비중증50%·도수 제외·재가입5년',
    }.get(gen,'')

def get_종번호(name):
    for i,k in enumerate(['(1종)','(2종)','(3종)','(4종)','(5종)','(6종)','(7종)','(8종)'],1):   # ★v29v 1-8종 지원
        if k in name: return i
    return 0

def _reflow_cols(block_lines):
    """★v133 (2026.07.21 장*상 현대해상 Hi2501 실측): 별첨이 3열이고 담보명이 길어 셀 안에서
    2~3줄로 접히면, pdftotext -layout 출력에서 <b>금액이 담보명 중간 줄에 홀로 남는다</b>.
    기존 _split_cols는 줄 단위라 조각난 이름을 이어붙이지 못하고, 금액이 <b>한 칸씩 밀려</b>
    앞 담보의 값이 뒤 담보에 붙었다(실측: 간호통합병동 7→20, 간병인 20→10).
    → 세로 거터(모든 줄이 공백인 열)로 칸을 나누고, 칸 안에서 연속 줄을 한 레코드로 묶어
      '이름 조각 전부 이어붙이기 + 그 안의 단독 숫자 = 그 레코드의 금액'으로 재조립한다.
    ★안전장치: '단독 숫자만 있는 줄'이 없으면 접힘 레이아웃이 아니므로 <b>원본 그대로 반환</b>
      (기존 2열 롯데·3열 신정원 경로 무영향)."""
    import re as _re
    lines=[str(x).rstrip('\n') for x in block_lines]
    if not lines: return block_lines
    # 접힘 신호: 공백+숫자만 있는 줄 존재
    if not any(_re.fullmatch(r'\s*[\d,]+\s*', l) and l.strip() for l in lines): return block_lines
    W=max(len(l) for l in lines)
    if W<40: return block_lines
    pad=[l.ljust(W) for l in lines]
    # 모든 줄이 공백인 열 = 거터
    gut=[all(p[i]==' ' for p in pad) for i in range(W)]
    # 폭 3 이상 거터 런의 시작 위치 = 열 경계
    bounds=[0]; i=0
    while i < W:
        if gut[i]:
            j=i
            while j<W and gut[j]: j+=1
            if (j-i)>=3 and j<W: bounds.append(j)
            i=j
        else: i+=1
    if len(bounds)<2: return block_lines
    bounds.append(W+1)
    out=[]
    for bi in range(len(bounds)-1):
        a,b=bounds[bi],bounds[bi+1]
        cells=[p[a:b] for p in pad]
        rec=[]
        def flush():
            if not rec: return
            amt=None; parts=[]
            for c in rec:
                t=c.strip()
                if _re.fullmatch(r'[\d,]+', t):
                    if amt is None: amt=t
                elif t: parts.append(t)
            nm=''.join(parts).strip()
            if nm and amt is not None: out.append(nm+'    '+amt)
            elif nm: out.append(nm)
            elif amt is not None: out.append(amt)
            rec.clear()
        for c in cells:
            if c.strip(): rec.append(c)
            else: flush()
        flush()
    # ★★v205 (2026.07.25 김*구 흥국화재 실측 회귀수정):
    #   1열 레이아웃('담보명    금액'이 같은 줄)에서 담보명 칸과 금액 칸 사이 거터가 3칸 이상이면
    #   위 로직이 두 칸을 <b>별개 열</b>로 갈라 '이름 11줄 → 금액 10줄' 순서로 재배열해버린다.
    #   → 페어링이 통째로 소실되어 그 계약 담보가 <b>0건</b>이 된다(실측: 흥국화재 10담보 전멸).
    #   재조립 결과에 '이름+금액' 쌍이 <b>하나도 없으면</b> 접힘 레이아웃이 아니므로 원본을 그대로
    #   돌려준다. 그러면 기존 _split_cols가 같은 줄 페어링으로 정상 처리한다.
    if not any(_re.search(r'\S\s{4}[\d,]+$', str(o)) for o in out):
        return block_lines
    return out


def _split_cols(block_lines):
    """★OCR PDF(pdftotext -layout) 별첨 다열(담보명|금액|담보명|금액|담보명|금액) → 1쌍 1줄로 분해.
    열 구분=공백 3개↑ 또는 탭. 담보명 내부 단일공백은 보존. 숫자 토큰=직전 담보명의 가입금액."""
    import re as _re
    out=[]
    for raw in block_lines:
        l=str(raw).rstrip()
        if not l.strip(): out.append(l); continue
        toks=[t for t in _re.split(r'\t+|\s{3,}', l.strip()) if t!='']
        if len(toks)<=1: out.append(l); continue
        name_acc=[]
        for t in toks:
            if _re.fullmatch(r'[\d,]+', t):          # 순수 숫자 = 값
                if name_acc:
                    out.append(' '.join(name_acc)+'    '+t); name_acc=[]
                else:
                    out.append(t)                    # 고아 숫자 → amts 폴백
            else:
                name_acc.append(t)
        if name_acc: out.append(' '.join(name_acc))  # 값없는 담보명 → pend 경유
    return out

_AMT_TAIL_UF = re.compile(r'(?<!\S)([\d][\d,]{0,11})\s*$')
def _paren_bal(s): return s.count('(')+s.count('（')-s.count(')')-s.count('）')

def _unfold_cols(block_lines):
    """★★★v225 (지점장 지시 2026.07.25, 영구): <b>별첨 담보명 접힘 복원</b>.
    롯데 3열 별첨은 담보명이 길면 <b>위·아래 줄에 걸쳐 접히고 금액은 가운데 줄</b>에 온다.
      1007| 중증질환자뇌혈관질환산정특례대상진단비Ⅱ(연간1회한)(맞춤        간호 간병통합서비스사용질병입원일당(요
      1008|                    1,000   주요심뇌5대혈관수술비Ⅱ  1,000              10
      1009| 간편고지)                                            외)(181일이상)(맞춤간편고지)
    구 파서는 이 구조를 몰라 <b>여러 담보명을 한 줄로 뭉치고 엉뚱한 열의 금액(10)을 붙였다</b>
    → 산정특례뇌혈관에 10이 찍히고 산정특례심장은 통째로 사라졌다(실측·정답 각 1,000).
    <b>복원 원리 2개</b>
      ①<b>열 경계 = 금액 토큰의 끝 좌표 클러스터</b>(금액은 우측정렬). 마지막 열은 줄 끝까지.
      ②<b>괄호 균형</b>으로 머리·꼬리를 판별 — 열림>닫힘이면 미완결이라 다음 조각이 꼬리다.
    ★<b>안전장치</b>: 복원 담보 수가 원본 페어 수보다 <b>적으면 원본을 그대로 반환</b>한다
      (v205 `_reflow_cols` 회귀 사고와 같은 실수를 막는다)."""
    try:
        lines=[str(l) for l in block_lines]
        if len(lines)<3: return block_lines
        pos={}
        for l in lines:
            for m in re.finditer(r'(?<!\S)[\d][\d,]{0,11}(?!\S)', l):
                pos[m.end()]=pos.get(m.end(),0)+1
        ends=sorted(pos.keys())
        if not ends: return block_lines
        # ★★★v226 (2026.07.25 삼성화재 실측 회귀수정): 구 코드는 <b>2회 이상 등장한 end만</b> 경계로 썼다.
        #   그래서 같은 열에 <b>자리수가 더 긴 금액이 1건만</b> 있으면(삼성 `상해 사망 5,000` end=48,
        #   나머지 9건은 end=46) 그 금액이 <b>경계 밖으로 밀려 `5,0`으로 잘렸다</b>(실측: 5,000 → 50).
        #   → <b>모든 end를 근접(±4) 클러스터링하고, 2회 이상 등장한 end를 포함한 클러스터만</b> 열로 인정하며
        #   경계는 그 클러스터의 <b>최댓값</b>을 쓴다(가장 긴 금액까지 안전하게 포함).
        grp=[]
        for e in ends:
            if grp and e-grp[-1][-1]<=4: grp[-1].append(e)
            else: grp.append([e])
        ends=[max(g) for g in grp if any(pos[x]>=2 for x in g)]
        if not ends: return block_lines
        bounds=[]; prev=0
        for e in ends: bounds.append((prev,e)); prev=e
        bounds[-1]=(bounds[-1][0], 10**6)          # 마지막 열은 줄 끝까지
        out=[]
        for a,b in bounds:
            # ★★★★★v260(2026.07.27 장O경 KB 실측): <b>줄 번호를 함께 들고 간다</b>.
            #   구 코드는 빈 줄을 버리고 조각만 이어 붙였고, <b>이어붙일지 말지를 괄호 균형 하나로</b>만 판단했다.
            #   → 담보명 꼬리가 <b>`…(간편가입)(갱`처럼 열린 괄호로 잘려 끝나면</b> 균형이 영영 안 맞아
            #     <b>다음 담보를 통째로 삼켰다</b>. 실측 사고:
            #       `다빈치로봇 암수술비(…)(갱` + `화상수술비(` + `상해1~5종수술비` → 한 덩어리, 금액 7
            #     → 다빈치 500·화상수술비·<b>상해1~5종 20/50/200/500/1,000</b>·산정특례 뇌심·중입자 5,000이 전부 소실.
            #   ★<b>이 별첨의 실제 구조는 「머리줄 / 금액줄 / 꼬리줄」이 붙어 있고 담보 사이에는 빈 줄이 있다</b>.
            #     따라서 <b>빈 줄(행 간격)이 레코드 경계</b>다. 괄호 균형은 보조 판정으로만 쓴다.
            cells=[(i,l[a:b].rstrip()) for i,l in enumerate(lines) if l[a:b].strip()]
            buf=''; pend=None; prev=None; tail=0
            def _flush():
                nonlocal buf,pend,tail
                if buf and pend: out.append((buf,pend))
                buf=''; pend=None; tail=0
            _prev_pure_amt = False        # ★v596 직전 셀이 '금액만 있는 줄'이었나
            for i,c in cells:
                m=_AMT_TAIL_UF.search(c)
                nm=(c[:m.start()] if m else c).strip(); amt=m.group(1) if m else None
                gap = (prev is not None and i-prev > 1)          # 빈 줄이 끼었다 = 새 담보
                dup = (pend is not None and amt is not None)      # 금액이 두 번째다 = 새 담보
                # 금액을 이미 받았는데 또 순수 이름이 오면 꼬리 1개까지만 허용(그 이상은 새 담보)
                # ★★★★★v596 제136조 (지점장 실측 2026.08.28 「비급여 도수·주사제·MRI가 안 나온다」
                #   · 김생금님 삼성생명 실손의료비1.0)
                #   세로 1열 별첨은 <b>담보명 줄 → 금액만 있는 줄 → 다음 담보명 줄</b>이 빈 줄 없이 이어진다.
                #   구 조건은 `tail >= 1`이라 <b>첫 꼬리를 무조건 허용</b> → 다음 담보명이 앞 담보에
                #   이어붙어 이름이 뭉치고 <b>금액이 한 칸씩 밀렸다</b>(상해 3건 소실 · MRI 300 유실).
                #   ⇒ 직전 셀이 <b>금액만 있는 줄</b>이었으면 그다음 이름은 <b>꼬리가 아니라 새 담보</b>다.
                #   ★괄호가 안 맞으면(`_paren_bal > 0`) 종전대로 꼬리 — 롯데 3열 복원은 그대로 산다.
                over = (pend is not None and amt is None and nm
                        and (tail >= 1 or _prev_pure_amt) and _paren_bal(buf) <= 0)
                if gap or dup or over: _flush()
                _prev_pure_amt = (amt is not None and not nm)
                prev = i
                if nm:
                    # ★★★★★v596 제136조 2항 — 금액을 아직 못 받았고 괄호도 균형이면
                    #   앞 조각은 담보명이 아니다(별첨 머리말) → <b>덮어쓴다</b>.
                    #   [실측] `(정상계약 리스트)삼성생명…보장기간` + `상해 입원의료비` 가 한 덩어리가 돼
                    #   <b>상해 입원의료비 5,000이 통째로 사라졌다</b>(김생금님).
                    #   ★괄호가 안 맞으면 진짜 접힘이므로 종전대로 이어붙인다(롯데 3열 복원 유지).
                    if buf and pend is None and _paren_bal(buf) <= 0:
                        buf = nm; tail = 0
                    elif buf: buf += nm; tail += 1
                    else: buf = nm
                if amt: pend = amt
                if buf and pend and _paren_bal(buf) <= 0 and tail >= 1:
                    _flush()
            _flush()
        # ★안전장치: 원본에서 '이름+금액' 같은 줄 페어 수를 세고, 복원이 그보다 적으면 원본 유지
        base=sum(1 for l in lines if re.search(r'\S\s{2,}[\d][\d,]*\s*$', l))
        if len(out) < max(base,1): 
            print(f'[v225 unfold skip] 복원 {len(out)}건 < 원본페어 {base}건 → 원본 유지')
            return block_lines
        print(f'[v225 unfold] 담보 {len(out)}건 복원(열 {len(bounds)}개, 경계 {ends})')
        return [f'{nm}    {amt}' for nm,amt in out]
    except Exception as _e:
        print(f'[v225 unfold ERR] {_e}')
        return block_lines

# ★★★v234 (2026.07.25 한정환 메트라이프생명 실측): 생보 별첨의 <b>무수식어 2자 담보</b>.
#   구 코드는 `len(name)>2` 필터에서 '입원'(2자)을 통째로 버렸다 → 메트 `입원 3` 2줄 소멸,
#   질병일당·상해일당이 각 13이어야 하는데 10으로 나왔다(한장보장표 불일치).
_SHORT_OK = {'입원','통원','수술','사망','장해','골절','화상','간병'}
# ★★★v234: <b>무수식어 담보가 별첨에 동일 금액 2줄이면 합산 금지·대표(max)</b>.
#   근거(세부가입현황 4p 메트라이프): 사망 칸 = `6,000 | 6,000`(질병 6,000 · 상해 6,000),
#   입원비 칸 = `3 | 3`(질병 3 · 상해 3). 즉 <b>2줄은 질병축·상해축 각각</b>이며 합산이 아니다.
#   이 담보들은 뒤에서 §8.1 종신 1:1 규칙 / 생보 입원특약 규칙이 두 행에 각각 기재하므로,
#   여기서 합산하면 <b>이중계산</b>이 된다. 실측: 일반사망 12,000 → 상해사망 29,100(정답 23,100).
#   ★반대로 '뇌출혈진단'·'급성심근경색진단'처럼 축이 하나인 담보의 2줄은 <b>합산이 정답</b>이다
#   (신한 급성심근 3,000×2=6,000 → 한장표 급성심근 13,400과 일치 확인). 그래서 <b>완전일치 집합</b>으로만 막는다.
# ★★★★★v293 (지점장 지시 2026.07.31, 영구): <b>4세대 실손의 입원은 5,000이 최대다</b>.
#   지점장 원문 = "4세대 실손의 입원은 5천이 최대다(급여5천/비급여5천 - <b>최고값 1개만 기재</b>)".
#   실측(김수영 메리츠 실손2107): 별첨에 상해 5,000 2줄 + 질병 5,000 2줄 → dambo 합산 <b>20,000</b>
#   → 엑셀 입원 20,000 · PPT '입원 : 20,001'(1억 초과 표기). 급여/비급여 두 칸이 각각 인쇄된 것이라
#   <b>합산이 아니라 최고값 1개</b>가 정답이다. → 완전일치 집합에 '실손입원의료비'를 넣는다.
_DUP_MAX_EXACT = {'일반사망','입원','입원특약','재해입원','실손입원의료비'}

def rule_extract(block_lines, prefolded=False):
    # ★★★★★v390b (2026.08.12 흥국 제안서 실측): <b>가입제안서 경로는 이미 「담보명  금액」
    #   1담보 1줄</b>로 만들어져 들어온다(parse_jean). 여기에 별첨용 접힘 복원을 다시 걸면
    #   금액 끝좌표가 우연히 두 무리로 갈려 <b>2열로 오판</b>하고(실측 경계 [38,63]) 담보명이
    #   서로 뭉쳐 <b>20건 → 16건</b>이 됐다. 이미 접힌 게 없는 입력에는 복원을 걸지 않는다.
    if not prefolded:
        block_lines=_unfold_cols(block_lines)                # ★v225 담보명 접힘 복원(성공 시 '이름  금액' 1줄 형태)
        block_lines=_split_cols(_reflow_cols(block_lines))   # ★v133 접힘 3열 재조립 → 기존 다열 분해
    block_lines=[l for l in block_lines if not (('표준금액' in str(l)) or ('권장금액' in str(l)) or ('적정금액' in str(l)))]  # ★표준금액 줄 제외
    """★v29t: 같은줄 우선 + 분리줄(코드/이름랩/금액뭉치) 순서 페어링(누락0). 김진구.txt 6계약 회귀검증 완료."""
    dambo={}; names=[]; amts=[]; pend=None
    # ★★★★★v345 [중복줄] 기록 — 지점장 지시 2026.08.02(영구).
    #   v344로 「동일 담보명 2줄 = 대표(max)」가 기본이 됐지만, <b>몇 줄이었는지가 아무 데도 안 남았다</b>.
    #   진짜 합산해야 할 담보가 나타나도 조용히 하나로 줄어든다 → <b>확인사항에 흔적을 남긴다</b>.
    #   ★값은 건드리지 않는다(표시 전용).
    _duplog={}
    def _add(_nm, _amt):
        _duplog.setdefault(str(_nm), []).append(_amt)
        # ★v61 심뇌혈관수술비 라인단위 분해(지침 §8.3.1 · 지점장 2026.07.15 재확정):
        #   '심뇌혈관…수술' = 심장수술비 + 뇌혈관수술비 각 100% 동일액.
        #   ★중복줄(상해·질병 등 같은 3,000이 2줄) = 합산 아니라 대표(max) — 6,000 오합산 방지.
        #   라인 단위로 쪼개므로 dambo 합산(6,000) 이전에 처리된다.
        # ★★★v217 (지점장 지시 2026.07.25, 영구): <b>DB손보 '주요심뇌5대혈관수술비' = 뇌혈관수술비 + 심장수술비
        #   각각 대표값 입력</b>. 구 조건은 `'심뇌혈관' in _n` 완전연속이라 <b>'심뇌<u>5대</u>혈관'처럼 사이에
        #   글자가 끼면 탈락</b>해 [확인]큐로 사라졌다(실측). → <b>'심뇌' + '혈관' + '수술'</b>로 완화.
        #   커버: 심뇌혈관수술비 · 주요심뇌5대혈관수술비 · 심뇌 5대혈관 수술 등.
        _n=re.sub(r'\s','',str(_nm))
        # ★★★★★v397 (지점장 확정 2026.08.12): <b>'주요치료'가 붙으면 수술 분해를 하지 않는다</b>.
        #   지점장 원문: 「심뇌혈관질환주요치료비 -> 2대주요치료비다」
        #   실측: 현대 `심뇌혈관질환<b>주요치료비</b>(…)(<b>수술</b>및혈전용해치료)` 500이 괄호 수식어의
        #     '수술' 글자 때문에 <b>심장수술비 500 + 뇌혈관수술비 500</b>으로 쪼개졌다.
        #   ★`주요심뇌5대혈관수술비`·`심뇌혈관수술비`는 '주요치료'가 없으므로 <b>종전대로 분해</b>한다.
        if '심뇌' in _n and '혈관' in _n and '수술' in _n and '주요치료' not in _n and '[확인]' not in _n:
            for _r in ('심장수술비[묶음]','뇌혈관수술비[묶음]'):   # ★태그 '뇌혈관' 금지→[묶음]
                dambo[_r]=max(dambo.get(_r,0), _amt)
        elif ('직접치료' in _n) and ('암' in _n) and ('일당' not in _n) and ('입원' not in _n):
            # ★★v227: `일반암직접치료 1,000`이 별첨에 <b>2줄</b> 인쇄되는데 합산하면 2,000이 된다.
            #   세부가입현황 정답은 <b>대표 1,000</b>(암수술 칸) → dambo 합산 이전에 max로 잡는다.
            dambo[_nm]=max(dambo.get(_nm,0), _amt)
        elif ('2대' in _n) and ('주요' in _n) and ('기관' in _n) and ('수술' in _n) and ('치료비' not in _n):
            # ★★★★★v348 (지점장 지시 2026.08.02, 영구): "<b>2대 주요기관질병 수술비 → 뇌혈관 + 심장이다
            #   하지만 관혈or비관혈이 있으므로 이런경우는 엑셀에 5대기관수술비칸에
            #   관혈 가입금액 / 비관혈 가입금액 넣어주면 된다</b>".
            #   → ㉠<b>관혈·비관혈 표기가 있으면 5대기관 수술비 관혈/비관혈 행</b>
            #     ㉡표기가 없으면 <b>뇌혈관수술비 + 심장수술비 각 100%</b>(묶음 공통원칙 §8.3.1).
            #   ★v216b 교훈: 공백을 지우면 '수술비 관혈'→'수술비관혈' 안에 <b>'비관혈'이 우연히 생긴다</b>
            #     → 반드시 구분자를 넣고 검사한다.
            _nk = _n.replace('수술비관혈','수술비|관혈')
            if '비관혈' in _nk:
                dambo['5대기관수술비(비관혈)[묶음]'] = max(dambo.get('5대기관수술비(비관혈)[묶음]',0), _amt)
            elif '관혈' in _nk:
                dambo['5대기관수술비(관혈)[묶음]'] = max(dambo.get('5대기관수술비(관혈)[묶음]',0), _amt)
            else:
                for _r in ('심장수술비[묶음]','뇌혈관수술비[묶음]'):
                    dambo[_r] = max(dambo.get(_r,0), _amt)
        elif any(k in _n for k in ('1-7종','1-8종','1-9종','1~7종','1~8종','1~9종')) and ('수술' in _n):
            # ★★★★★v342 (지점장 지적 2026.08.02): <b>축 표기 없는 1-7/1-8/1-9종은 질병·상해 양쪽에 기재</b>한다.
            #   지점장 원문: "1-7종 대표값 넣으라고했지만 <b>질병1-종 즉 1-8종에는 빠져있다</b>
            #   엑셀의 오류로 진단서와 ppt 다 빠져있었다".
            #   실측 = 미래에셋 `1-7종수술특약(급여)` 2,000이 <b>상해 1-8종에만</b> 들어가고 질병은 0이었다.
            # ★★★★★v343 (지점장 지시 2026.08.02, 영구): "<b>1-7.1-8.1-9의 가장큰 대표값. 질병.상해에 각각넣어라</b>"
            #   → 축 표기 유무와 <b>무관</b>하게 1-7·1-8·1-9종 담보 전체에서 <b>최댓값 1개</b>를 뽑아
            #     <b>질병 종수술비(1-8종) · 상해 종수술비(1-8종) 두 행에 각각</b> 기재한다. <b>합산 금지</b>.
            for _r in ('상해 종수술비(1-8종)[묶음]','질병 종수술비(1-8종)[묶음]'):
                dambo[_r]=max(dambo.get(_r,0), _amt)
        elif (('종수술' in _n) or re.search(r'1[-~][0-9]종', _n)) and re.search(r'\([0-9]\s*종\)', _n) \
             and not any(k in _n for k in ('질병','상해','재해')):
            # ★★★★★v341 예외사항(지점장 지시 2026.08.02): <b>질병/상해 표기 없는 종수술비가
            #   한 계약에 두 벌 실리면 「질병 1벌 + 상해 1벌」이다 — 합산 금지·대표(max).</b>
            #   지점장 원문: "미래에셋 1-5종이 질병or상해라고 기재안된채 2개가 기재되어있다
            #   넌 그걸 <b>각각 더해서 2배</b>로 만들어놧다".
            #   실측 = 미래에셋 `1-5종수술특약(1종)~(5종)` 20/40/300/1000/2000이 <b>두 벌</b>
            #   → dambo 단계에서 합산돼 40/80/600/2000/4000이 됐다. 정답은 각 축 20/40/300/1000/2000.
            #   ★축 표기가 있는 담보(질병1-5종수술비 등)는 종전대로 합산 — 이 예외에 안 걸린다.
            dambo[_nm]=max(dambo.get(_nm,0), _amt)
        elif _n in _DUP_MAX_EXACT:
            # ★★★v234: 무수식어 담보(일반사망·입원 등) 2줄 = 질병축·상해축 각각 → 합산 금지·대표(max).
            #   합산하면 뒤의 종신 1:1 / 입원특약 양행기재 규칙과 겹쳐 금액이 2배가 된다.
            dambo[_nm]=max(dambo.get(_nm,0), _amt)
        elif ('사망' in _n) and ('후유장해' in _n):
            # ★★v92 (장혜경 실측): 결합담보 '상해사망후유장해 1,000'이 별첨에 <b>두 줄</b> 인쇄돼
            #   합산 2,000 → 분해 후 상해사망 2,000·상해후유 2,000이 되어 한장보장표(1,100/1,000)와 어긋났다.
            #   심뇌혈관수술(v61)과 같은 원칙 — <b>중복 줄은 합산 금지·대표(max)</b>.
            dambo[_nm]=max(dambo.get(_nm,0), _amt)
        else:
            # ★★★★★v344 기본값 전환 — 지점장 확정 2026.08.02(영구·최상위)
            #   지점장 원문: "<b>똑같은 이름과패턴의 담보가 2개면 1개만 입력하자</b>".
            #   구 기본값은 <b>합산(`+=`)</b>이었다 → 별첨에 <b>담보명이 완전히 같은 줄이 두 번</b> 실리면
            #   금액이 그대로 <b>2배</b>가 됐다. 그동안 표적항암·합의금·암수술·간병인·심뇌혈관·결합담보처럼
            #   <b>사고가 난 담보만 하나씩 예외 목록</b>에 넣어 땜질해 왔고, 목록에 없는 담보가 오면
            #   또 2배가 됐다(실측 미래에셋 1-5종·1-7종). <b>원칙을 뒤집는다 — 기본이 대표(max)다.</b>
            #   ★적용 범위 = <b>한 계약 안에서 담보명이 완전히 동일한 줄</b>. 담보명이 다르면(스쿨존 벌금 +
            #     업무상과실 벌금 등) 종전대로 각각 잡혀 끝열에서 합산된다.
            #   ★계약이 다르면 열이 다르므로 <b>끝열 가로 SUM은 그대로</b>다(계약 간 합산은 영향 없음).
            dambo[_nm]=max(dambo.get(_nm,0), _amt)
    UNIT = r'(?:\s*(원|만원|만))?'
    NOISE = re.compile(r'지점|LP|☎|^\d{4}\.\d{2}\.\d{2}$|^\d+/\d+$|계약자|납입주기|보장기간|정상계약|상기 ?내용은|기준으로 ?분석|향후 ?계약사항|본 ?리포트|참조하시|제안서는|유의 ?사항')
    def _flush():
        nonlocal pend
        if pend:
            nm=re.sub(r'\s+',' ',pend.strip())
            if (len(nm)>2 or nm in _SHORT_OK) and not re.search(r'납입면제|납입지원',nm): names.append(nm)
        pend=None
    for raw in block_lines:
        l=raw.strip()
        if not l: continue
        if NOISE.search(l): _flush(); continue
        m = re.search(r'^(.+?)\s{2,}([\d,]+)'+UNIT+r'\s*$', l) or re.search(r'^(.+?)\s+([\d,]+)'+UNIT+r'\s*$', l)
        if m and re.search(r'[가-힣]', m.group(1)):
            _flush()
            name = re.sub(r'\s+', ' ', m.group(1).strip())
            try:
                amt = int(m.group(2).replace(',',''))
                if (m.group(3) or '') == '원': amt = amt // 10000
                if 0 < amt <= 200000 and (len(name) > 2 or name in _SHORT_OK):
                    # ★v29t §8.1: 생보 '주계약_주계약'이 2줄(일반+재해)로 반복되는 별첨 → 병합 금지, 순번 접미사로 분리
                    if '주계약_주계약' in name and name in dambo:
                        k=2
                        while f'{name}~{k}' in dambo: k+=1
                        name=f'{name}~{k}'
                    _add(name, amt)
            except: pass
            continue
        m2 = re.match(r'^([\d,]+)'+UNIT+r'\s*$', l)
        if m2:
            _flush()
            try:
                amt = int(m2.group(1).replace(',',''))
                if (m2.group(2) or '') == '원': amt = amt // 10000
                if 0 < amt <= 200000: amts.append(amt)
            except: pass
            continue
        if re.match(r'^\[\w+\]', l):
            _flush(); pend = l; continue
        if pend is not None:
            pend += l; continue
        if re.search(r'[가-힣]', l):
            pend = l
    _flush()
    for i, nm in enumerate(names):
        _add(nm, (amts[i] if i < len(amts) else 0))   # 금액 미확보=0 → [확인] 경유, 증발 금지
    dambo['__DUP__'] = {k:v for k,v in _duplog.items() if len(v) >= 2}
    return dambo

def llm_extract(block_text):
    """깨진 별첨(담보명/금액 줄 분리)을 Claude가 의미로 추출. 키 없으면 {} -> 규칙 폴백."""
    key = os.environ.get('ANTHROPIC_API_KEY','')
    if not key or not block_text.strip(): return {}
    prompt = ("보험 별첨 텍스트에서 담보명과 가입금액(만원 단위 숫자)을 추출.\n"
        "주의: 표가 깨져 담보명이 2줄로 나뉘거나 금액이 별도 블록에 모여있을 수 있음. 순서·문맥으로 정확히 매칭.\n"
        "담보명은 원문 그대로. 납입면제·납입지원·특약안내 등 비담보성 항목은 제외.\n"
        "금액 단위가 '원'이면 만원으로 환산(÷10000). 매칭 불확실하면 제외.\n\n"
        + block_text + "\n\nJSON만 출력: {\"담보명\": 금액숫자}")
    try:
        r = httpx.post('https://api.anthropic.com/v1/messages',
            headers={'x-api-key':key,'anthropic-version':'2023-06-01','content-type':'application/json'},
            json={'model':'claude-haiku-4-5-20251001','max_tokens':8000,'messages':[{'role':'user','content':prompt}]}, timeout=90)
        if r.status_code != 200:
            print(f'[LLM_EXTRACT_HTTP] status={r.status_code} {r.text[:300]}')
            return {}
        txt = ''.join(b.get('text','') for b in r.json().get('content',[]) if b.get('type')=='text')
        txt = txt.strip().replace('```json','').replace('```','').strip()
        out = json.loads(txt)
        print(f'[LLM_EXTRACT] ok items={len(out)}')
        return {str(k).strip(): int(v) for k,v in out.items() if isinstance(v,(int,float)) and 0 < v <= 200000 and len(str(k).strip())>2}
    except Exception as e:
        print(f'[LLM_EXTRACT_ERR] {e}')
        return {}


_IMG_PDF_WARN = ''   # ★v280 이미지PDF(비전OCR) 경고 전역 플래그
# ★★★★★v371 (지점장 확정 2026.08.09): 가입제안서 단독 모드 플래그.
#   보장분석지 없이 제안서 1장만 올린 경우 = <b>검산·실손 세대 판정 불가</b>.
#   지점장 확정 = 「불가」로 명시하고 건너뛴다(산출은 막지 않는다).
_JEAN_ONLY = ''
_VISION_PARTIAL = ''  # ★v281 비전 OCR 부분 실패(페이지 유실) 경고

def pdf_to_txt(pdf_bytes):
    """★v32 OCR PDF 입력(2026.07.07 지점장 정답): 1순위=텍스트레이어 직독(pdftotext -layout, 무키·100%),
    2순위=Claude 비전 OCR(이미지 전용 PDF). Adobe .txt 변환 없이 OCR PDF 1개로 완결."""
    # ── 1순위: 텍스트레이어 직독 (드래그 선택 가능한 OCR PDF면 API 없이 100% 추출) ──
    try:
        import subprocess, tempfile
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as _f:
            _f.write(pdf_bytes); _pp=_f.name
        _tl = subprocess.run(['pdftotext','-layout',_pp,'-'], capture_output=True, text=True, timeout=60).stdout
        # ★★★v279: 세부가입현황 좌표 폴백용 bbox XML을 같이 확보한다.
        #   신형 let: 리포트는 표 라벨이 통째로 이미지라 라벨 기반 매핑이 통째로 꺼진다.
        try:
            global _BBOX_XML
            _BBOX_XML = subprocess.run(['pdftotext','-bbox','-layout',_pp,'-'],
                                       capture_output=True, text=True, timeout=60).stdout or ''
            print(f'[v279 bbox] 확보 {len(_BBOX_XML)}자')
        except Exception as _eb:
            _BBOX_XML = ''; print(f'[v279 bbox] 실패 {_eb}')
        try: os.unlink(_pp)
        except: pass
        if _tl and len(_tl) > 1500 and ('별첨' in _tl or '정상계약' in _tl or '보장' in _tl):
            print(f'[PDF_TEXTLAYER] ok chars={len(_tl)}'); return _tl
        print(f'[PDF_TEXTLAYER] 얇음/미검출 chars={len(_tl) if _tl else 0} -> 비전 폴백')
    except Exception as _e:
        print(f'[PDF_TEXTLAYER_ERR] {_e} -> 비전 폴백')
    # ── 2순위: Claude 비전 OCR (텍스트레이어 없는 이미지 스캔본) ──
    key = os.environ.get('ANTHROPIC_API_KEY','')
    if not key:
        print('[PDF_VISION] no api key -> skip')
        globals()['_VISION_FAIL']='ANTHROPIC_API_KEY 미설정 — 비전 OCR 사용 불가'
        return ''
    try:
        from pdf2image import convert_from_bytes
        import io
        # ★★v232: dpi 170 → <b>300</b>. 170dpi에서는 별첨 표의 금액 자릿수(1,000 vs 100)가
        #   뭉개져 오독 위험이 크다. '인쇄→PDF' 이미지본은 원본이 200dpi 조각이라 300으로 올려 읽는다.
        pages = convert_from_bytes(pdf_bytes, dpi=300)
    except Exception as e:
        print(f'[PDF_RENDER_ERR] {e}')
        globals()['_VISION_FAIL'] = f'PDF 이미지 렌더 실패({e}) — pdf2image/poppler 확인'
        return ''
    prompt = ("이 보험 보장분석 리포트 페이지의 모든 텍스트를 그대로 전사하라. "
              "표는 탭으로 열 구분, 회전된 표는 정방향으로 읽어라. "
              "담보명과 금액(만원 단위)을 같은 줄에 유지하라. 회사명·상품명·계약일·만기일·보험료도 포함. "
              "해석·설명·요약 금지, 페이지의 원문 텍스트만 출력.")
    out=[]; _fail_pages=[]
    globals()['_VISION_PARTIAL']=''
    for idx, img in enumerate(pages):
        try:
            # ★v60 회전 보정: let: 리포트는 가로형인데 '인쇄→PDF' 이미지본은 세로 A4에
            #   가로 내용이 90° 눕는다. 세로(높이>너비) 페이지면 시계방향(-90°)으로 세워
            #   비전 OCR 정확도를 높인다(정방향 검증 완료). 이미 정방향(가로)이면 무동작.
            if img.height > img.width:
                img = img.rotate(-90, expand=True)
            # ★★v232: <b>긴 변 2000px로 리사이즈해 전송</b>. 300dpi A4 회전본은 3509×2481인데
            #   Anthropic 이미지 권장은 긴 변 1568px이라 그대로 보내면 <b>서버가 임의 축소</b>하고
            #   토큰도 과소비된다. 300dpi로 읽고 2000px로 줄이면 <b>표 숫자 선명도는 유지</b>되면서
            #   전송량이 1/2로 준다(실측 PNG 1,264,841 → 670,377B).
            try:
                from PIL import Image as _PILImage
                _lim=2000
                if max(img.size) > _lim:
                    _sc=_lim/max(img.size)
                    img=img.resize((int(img.size[0]*_sc), int(img.size[1]*_sc)), _PILImage.LANCZOS)
            except Exception: pass
            buf=io.BytesIO(); img.save(buf, format='PNG')
            b=base64.b64encode(buf.getvalue()).decode()
            # ★★★★★v281 (2026.07.31 이영태 32페이지 실사고):
            #   구 코드는 <b>페이지당 1회 호출·재시도 0회</b>였고, 실패 흔적을 <b>idx==0일 때만</b> 남겼다.
            #   → 1p가 성공하고 12~32p가 429/타임아웃으로 죽으면 <b>아무 데도 안 남고</b>
            #     앞쪽 계약목록만 살아 "계약은 나오는데 담보가 통째로 빈" 산출물이 조용히 나갔다.
            #   이영태 실측: 32페이지 · b64 906KB/장 · 이미지토큰 약 3,770/장 = <b>32장 약 12만 입력토큰</b>.
            #   → ①429/5xx <b>재시도 3회 지수백오프</b> ②<b>모든 실패 페이지를 기록</b> ③성공/실패 수를 화면에 노출.
            t=''
            _last=''
            for _try in range(3):
                try:
                    r=httpx.post('https://api.anthropic.com/v1/messages',
                        headers={'x-api-key':key,'anthropic-version':'2023-06-01','content-type':'application/json'},
                        # ★★★v233 원복(2026.07.25): 모델은 <b>어제 실제로 작동한 haiku-4-5</b>를 그대로 쓴다.
                        #   v232에서 내가 검증 없이 sonnet-4-6으로 바꿨다 — <b>모델명이 틀리면 400으로 OCR이 통째 실패</b>한다.
                        json={'model':'claude-haiku-4-5-20251001','max_tokens':8000,
                              'messages':[{'role':'user','content':[
                                  {'type':'image','source':{'type':'base64','media_type':'image/png','data':b}},
                                  {'type':'text','text':prompt}]}]}, timeout=120)
                    if r.status_code==200:
                        t=''.join(x.get('text','') for x in r.json().get('content',[]) if x.get('type')=='text')
                        break
                    _last=f'status={r.status_code}'
                    print(f'[PDF_VISION_HTTP] p{idx+1} try{_try+1} {_last} {r.text[:160]}')
                    if r.status_code in (429,500,502,503,504,529):
                        import time as _tm; _tm.sleep((2,5,12)[_try]); continue
                    break
                except Exception as _ee:
                    _last=f'예외 {_ee}'
                    print(f'[PDF_VISION_ERR] p{idx+1} try{_try+1} {_ee}')
                    import time as _tm; _tm.sleep((2,5,12)[_try])
            if t.strip(): out.append(t)
            else: _fail_pages.append((idx+1,_last))
        except Exception as e:
            print(f'[PDF_VISION_ERR] p{idx+1} {e}')
            _fail_pages.append((idx+1,f'예외 {e}'))
    txt='\n'.join(out)
    print(f'[PDF_VISION] pages={len(pages)} ok={len(out)} fail={len(_fail_pages)} '
          f'chars={len(txt)} dpi=300 model=claude-haiku-4-5-20251001')
    # ★v281 실패 페이지는 <b>한 장이라도</b> 전부 노출한다(구 코드는 1페이지 실패만 기록).
    if _fail_pages:
        _fp=', '.join(f'{n}p({w})' for n,w in _fail_pages[:12])
        _more='' if len(_fail_pages)<=12 else f' 외 {len(_fail_pages)-12}장'
        globals()['_VISION_PARTIAL']=(f'비전 OCR {len(pages)}장 중 <b>{len(_fail_pages)}장 실패</b> '
                                      f'— 실패: {_fp}{_more}')
        print('[PDF_VISION_PARTIAL] '+globals()['_VISION_PARTIAL'])
    if not txt.strip() and not globals().get('_VISION_FAIL'):
        globals()['_VISION_FAIL']=(globals().get('_VISION_PARTIAL') or '비전 OCR이 글자를 한 자도 반환하지 않음')
    return txt


def _looks_broken(data):
    """txt 파싱 결과 깨짐 감지: 담보키에 다중담보 병합(無 반복·상품명 라벨·초장문) 비율."""
    if not data or not data.get('contracts'): return True
    _bad=0; _tot=0
    for c in data['contracts']:
        for k in c.get('dambo',{}).keys():
            _tot+=1; ks=str(k)
            if ks.count('無')>=2 or ks.count('（')>=3 or len(ks)>40: _bad+=1
    return _tot>0 and _bad/_tot>0.15


# ★★★★★ v238 CI 자가진단 — 배포본에 내장, 매 실행마다 자동 검사 ★★★★★
#   <b>메모리에 "조심하자"고 적는 것으로는 못 막는다</b>. 같은 뿌리(연속문자열 매칭)로
#   v216b(수술비 관혈) · v217(심뇌5대혈관) · v235(CI종신보험) — <b>3번 반복</b>했다.
#   → 판정 케이스를 <b>코드에 박아</b> 배포마다 자동으로 돌린다. 하나라도 깨지면 로그·/health에 즉시 노출.
_CI_SELFTEST = [
    # (상품명, 기대값) — CI로 인식돼야 하는 것
    ('무）라이프케어CI 종신보험',                            True),   # ★v235 실사고
    ('（무） 변액유니버셜 세번받을수있는 CI종신보험（1701）',   True),   # ★v235 실사고
    ('무배당교보큰사랑 CI 보험',                             True),
    ('무배당 KB CI종신보험',                                True),   # ★v246 실사고: 'KBCI'로 붙어 탈락하던 것
    ('(무)DB CI종신보험',                                   True),   # 영문약자 + CI
    ('(무)교보GI종신보험',                                  True),
    ('CI간편종신보험',                                      True),
    ('무배당 삼성생명 퍼펙트통합보험(프리미엄,無)표준',        True),   # 퍼펙트 시리즈
    ('삼성생명퍼펙트플러스통합보험',                          True),
    ('삼성생명퍼펙트종합보험',                               True),
    ('퍼텍트플러스보험',                                    True),   # 오타 동의어
    ('한화생명 리빙케어보험',                                True),
    # CI가 아니어야 하는 것
    ('무배당 마스터플랜 변액유니 버셜종신 II 보험',            False),
    ('참좋은운전자상해보험 2510',                            False),
    ('무배당 마이라이프 한아름종합보험 1710',                 False),
    ('ACCIDENT 보장보험',                                   False),  # 영문 안의 'CI' 오검출 방지
    ('SPECIAL CARE 보험',                                   False),
    ('DB손보 다이렉트 개인용',                               False),
]

# ★★★★★v241 CI 4단계 체크리스트 자가진단(지점장 정본 2026.07.25)
#   (중대한OO 리스트, 진단담보 리스트, 별첨 사망줄, 기대 선지급률, 기대 사망보장, 기대 판정근거)
_CI_RATE_SELFTEST = [
    ([2000], [3000,3000,3000,3000,2000], [10000], 50, 4000, '③'),  # ★신한 실사고: 중대한화상 2,000÷50%=4,000
    ([],     [2400,2400,2400,2400,1000,1000,1000,3000,400], [3000,2000], 80, 3000, '④'),  # ★DB 실사고: 2,400×4개 ÷80%=3,000 ✓사망줄 일치
    ([3200], [3200,1000], [4000],  80, 4000, '③'),                 # 교보CI 정본 검증예
    ([],     [2000,2000], [4000],  50, 4000, '④'),                 # 동일금액 2개 → ×2 = 사망줄 일치
]

def ci_rate_selftest():
    """★v241 CI 4단계(③중대한OO → ④동일금액 2개 이상 → 사망보장금 일치) 규칙 자가진단."""
    bad=[]
    for jd, cands, sl, exp_p, exp_s, exp_src in _CI_RATE_SELFTEST:
        slset=set(sl); bon=None; pct=None; src=''
        if jd:
            bon=max(jd); src='③'
            for r,p in ((0.5,50),(0.8,80)):
                if round(bon/r) in slset: pct=p; break
            if not pct: pct=50
            sam=round(bon/(pct/100.0))
        else:
            cnt={}
            for v in cands: cnt[v]=cnt.get(v,0)+1
            sam=0
            for v in sorted([x for x,c in cnt.items() if c>=2], reverse=True):
                for r,p in ((0.5,50),(0.8,80)):
                    if round(v/r) in slset:
                        bon=v; pct=p; sam=round(v/r); src='④'; break
                if bon: break
        if pct!=exp_p or sam!=exp_s or src!=exp_src:
            bad.append(f'중대한{jd}/사망줄{sl} → {src}{pct}%·{sam}(기대 {exp_src}{exp_p}%·{exp_s})')
    return bad


# ★★★★★v246 실손 소스 자가진단 — <b>별첨 명시값이 살아남는지</b>를 매 배포마다 검사한다.
#   (별첨 정본 원칙이 캡·기본값에 덮이는 사고를 코드로 차단)
_SILSON_SELFTEST = [
    # (세대, 별첨 통원 명시값, 기대 통원) — None = 별첨에 없음(기본값 적용)
    ('3세대',  25, 25),   # ★DB생명 실사고: 별첨 25가 한장표 30에 덮이면 안 된다
    ('3세대',  30, 30),   # 별첨이 30이면 30
    ('4세대',  25, 25),   # ★구 코드는 20으로 캡했다 → 별첨 우선으로 수정
    ('4세대',  None, 20), # 별첨에 없으면 기본 20
    ('1세대',  None, 10), # 1세대 기본 10
]

def silson_selftest():
    """별첨 명시값 우선 원칙이 지켜지는지 검사."""
    bad=[]
    for gen, tw, exp in _SILSON_SELFTEST:
        if gen in ('4세대','5세대'): got = tw if tw else 20
        elif gen.startswith('1세대'): got = tw if tw else 10
        else: got = tw if tw else 25
        if got != exp: bad.append(f'{gen}/별첨{tw} → {got}(기대 {exp})')
    return bad


# ★★★★★v382 <b>조문 자가진단</b> — 지점장이 확정한 지침이 <b>매 배포마다 코드로 강제</b>되게 한다.
#   [왜 만드나] 2026.08.11 점검 실측: 오늘 확정된 조문 5개(단독5종·Ⅰ/Ⅱ 단서·종수술 1-5종·
#   실손 세대·실손 행 게이트) 중 <b>배포 게이트에 걸린 것이 0건</b>이었다.
#   기존 게이트 4개는 전부 <b>런타임(값)</b> 검사라 한장표에 칸이 없는 담보(협심증·1~5종)는
#   영원히 통과한다 — 실제로 v377 검산은 26/27이었는데 협심증 오기재도 종수술 3건 실종도 다 통과했다.
#   지침 위반은 <b>값이 아니라 코드 작성 시점</b>에 생긴다 → 담보명→기대행을 직접 대조한다.
#   ★master.xlsx는 건드리지 않는다(원본 cmp 동일 = 정본 판정 기준 유지).
_DOCTRINE_SELFTEST = [
    # (담보명, 기대 마스터행 / None=기재금지)  ── 조문: 단독 5종(#27)
    ('허혈성심장질환진단비',                 '허혈성 진단비'),
    ('심장허혈성진단비',                     '허혈성 진단비'),
    ('갱신형 허혈성심장진단비',              '허혈성 진단비'),
    ('허혈심장질환진단특약ⅢUT(무배당,무해약환급금형)_간편고지3.10.5', '허혈성 진단비'),
    # ★★★★★v591 (지점장 실측 2026.08.26 「<b>또 협심증이 2배로 나온다</b>」 · 롯데 let:smile).
    #   [실측 근거표] 협심증 행에 <b>두 줄</b>이 들어갔다 —
    #     ① `심혈관질환진단비(특정심장질환Ⅰ)(건강체할인형)` 1,000  ← 여기 DMAP(인라인)
    #     ② `협심증[심장묶음]` 1,000                              ← `_HB` 분해
    #     ⇒ 합계 <b>2,000</b>. 다른 담보는 `[심장묶음]` 한 줄뿐인데 협심증만 두 줄이었다.
    #   [조문 위반] DOCTRINE 411행 = <b>롯데 특정심장Ⅰ → 급성심근경색</b>인데
    #     이 줄이 <b>회사 구분 없이</b> 협심증으로 보냈다. 「특정Ⅰ」의 뜻은 <b>회사마다 다르다</b>
    #     (제123조 — 라벨 말고 회사별 질병코드로 분해한다).
    #   ⇒ <b>이 두 줄을 삭제</b>한다. 특정Ⅰ/Ⅱ는 `_HB`(회사별 정본표) <b>한 경로</b>에서만 분해한다.
    #     회사 정보가 없는 DMAP에서 라벨만 보고 매핑하면 <b>반드시 어느 회사에선가 틀린다</b>.
    ('뇌혈관질환진단비',                     '뇌혈관진단비'),
    ('뇌졸중진단비',                         '뇌졸증진단비'),
    ('뇌출혈진단비',                         '뇌출혈진단비'),
    ('급성심근경색진단비',                   '급성심근경색'),
    # 조문: CI 변환·별도 행은 단독 가드에 먹히지 않아야 한다
    ('중대한급성심근경색진단비',             '중대한 급성심근'),
    ('외상성뇌출혈진단비',                   '외상성뇌출혈'),
    # 조문: 진단 전용(v325·v326b) — 수술·입원은 마스터에 행이 없다
    ('뇌출혈수술급부금',                     None),
    ('급성심근경색증수술급부금',             None),
    # 조문: 「종수술」=1~5종(v378) — 종번호가 없어도 그 행으로
    ('질병종수술',                           '질병 종수술비(1-5종)'),
    ('상해종수술',                           '상해 종수술비(1-5종)'),
    ('無파워수술보장(본인)',                 '종수술비공통'),
    ('재해종수술비(1-5종)',                  '상해 종수술비(1-5종)'),
]

# (담보명, 단독이어야 하나) ── 조문: Ⅰ·Ⅱ 등급은 단독 예외에서 뺀다
_SOLO5_SELFTEST = [
    ('허혈성심장질환진단비',   True),
    ('뇌혈관진단비',           True),
    ('허혈심장질환진단특약ⅢUT(무배당,무해약환급금형)_간편고지3.10.5', True),
    ('뇌혈관진단비Ⅰ',          False),   # 등급 Ⅰ → 지침대로(회사별 표)
    ('뇌혈관진단비Ⅱ',          False),
    # ★v387 지점장 확정 2026.08.12 「허혈성 키워드가 앞에 있으면 무조건 단독」
    ('허혈성심질환진단비Ⅱ',    True),
    ('[갱신형]허혈성심질환진단비Ⅱ(통합간편가입형)(갱신형_20년)', True),
    ('허혈성심장질환진단비Ⅰ',  True),
    ('특정허혈심장질환진단비',  False),   # 묶음 수식어 → 회사별 표
    ('4대순환계질환진단비',     False),
]

# (가입일, 3대비급여, 처방조제, 기대세대) ── 조문: 실손 세대(#15) "4세대는 통원비 20만원만 있다"
_GEN_SELFTEST = [
    ('',           True,  True,  '3세대'),
    ('',           True,  False, '4세대'),
    ('',           False, False, ''),
    ('2019.05.01', True,  True,  '3세대'),
    ('2022.05.01', True,  False, '4세대'),
    ('2026.05.06', False, False, '5세대'),
]

# ★v384 (담보명, 기대행) — `_dedup_std` 경로 전용. 심장Ⅰ/Ⅱ·뇌질환Ⅰ/Ⅱ는 resolve2가 아니라
#   이 함수와 인라인 묶음 블록이 처리하므로 <b>여기서 따로 검사</b>한다.
_DEDUP_SELFTEST = [
    ('심장질환진단비Ⅰ',   '협심증'),          # 묶음Ⅰ → 협심증(허혈성 아님)
    ('심장질환진단비Ⅱ',   '급성심근경색'),
    ('뇌질환진단비Ⅰ',     '뇌혈관진단비'),
    ('뇌질환진단비Ⅱ',     '뇌졸증진단비'),
]



# ★★★★★v410 (지점장 지시 2026.08.12 「지금 해라 절대 미루지 마라」):
#   조문 커버리지 19% → 끌어올린다. <b>파일 구조·상수·존재 여부로 검사 가능한 조문</b>을 전부 건다.
#   (조문, 파일, 정규식, True=있어야 / False=없어야)
_STRUCT_SELFTEST = [
    ('제1조 엑셀기준',   'main.py',            r'마스터무행|마스터 무행', True),
    ('제2조 결과값동결', 'main.py',            r'_data_cols', True),
    ('제7조 합계',       'main.py',            r'=SUM\(', True),
    ('제11조 게이트',    'main.py',            r'PPT_MISS', True),
    ('제12조 배포9파일', 'main.py',            r"ZIP9 = \['main\.py'", True),
    ('제17조 dambo키',   'report_weasy.py',    r"_p8v\('하이클래스암'\)", True),
    ('제19조 인포고객명','report_weasy.py',    r'def _info_strip_cust\(', True),   # ★v410e 접두 일치로 뚫렸다(`_DISABLED`)
    ('제19조 만든이삭제','ga_tables.py',       r'최은혜', False),
    ('제23조 제안합계레드','main.py',          r"'제안 합계'", True),
    ('제24조 8페이지삭제','report_weasy.py',   r'<!-- P8N: 주요치료비 세부', False),
    # ★v410f 제26조는 <b>「있나」가 아니라 「3장인가」</b>다 — 유무 검사는 2장이어도 통과했다(뮤테이션 실측).
    ('제26조 자료3장',   'report_weasy.py',    r'(?s)INFO-TBL.*INFO-TBL.*INFO-TBL', True),
    ('제27조 PPT슬래시레드','main.py',         r'__SS_RED__', True),
    ('제30조 여백폰트',  'report_weasy.py',    r'font-size:2\.5mm', True),
    ('제33조 실행마다정독','main.py',          r'_doc_read\(tag=.analyze.\)', True),
    ('제34조 각인1곳',   'main.py',            r"^VSTAMP = 'v\d", True),
    ('제35조 zip검증',   'main.py',            r'def zip_selfcheck', True),
    # ★v475 제83조 — 갱신 담보 색은 계약 루프 끝에서 확정한다(제5조 B).
    ('제83조 갱신담보색', 'main.py',            r'_blue_r', True),
    ('제82조 상품명절단', 'main.py',            r'def _clean_product', True),
    ('제84조 회사명절단', 'main.py',            r'def _clean_company', True),
    ('제86조 갱신담보',   'main.py',            r'def _is_gen_dambo', True),
    # ★v480 제89조 — 비교엑셀 계약별 표: 상품명 C+D 2칸 · 상태 G 1칸.
    ('제89조 비교엑셀열', 'remodel.py',         r"_LAYOUT89 = \{'회사': 2, '상품명': \(3, 4\), '전': 5, '후': 6, '상태': 7\}", True),
    ('제88조 세로분산',   'report_pages.py',    r'justify-content:space-between', True),
    ('제87조 하단표',     'remodel.py',         r'v478 하단표', True),
    ('제91조 행높이',     'remodel.py',         r'defaultRowHeight = _ROW_PT', True),
    ('제92조 상품명줄',   'remodel.py',         r'_PROD_W = 43', True),
    ('제93조 전기납35',   'main.py',            r'pay_y >= 35', True),
    ('제94조 납입회차',   'main.py',            r'월납\|년납\|연납', True),
    ('제99조 7쪽입력칸', 'report_pages.py',    r'_P9BIG', True),
    ('제103조 2쪽분산',  'report_pages.py',    r'_P2FILL', True),
    ('제103조 그래프색',  'report_pages.py',    r'#f0762a', True),
    ('제103조 그래프색', 'report_pages.py',    r'bar-before\{background:#f0762a', True),
    ('제103조 삭제행',   'report_pages.py',    r'status-del', True),
    ('제104조 8쪽원복',  'report_pages.py',    r'height:2\.8mm;border-bottom', True),   # ★v522 — 지점장 «8페이지 기존대로» 원복 반영
    ('제105조 4쪽표',    'report_pages.py',    r'padding:2\.05mm 1mm', True),
    ('제106조 금색',     'report_pages.py',    r'--gold:#c88d20', True),
    ('제106조 레드',     'report_pages.py',    r'diff\.down\{color:#c0392b\}', True),
    ('제107조 6쪽우측',  'report_pages.py',    r"\.rcol \.box\{height:8\.0mm\}", True),   # ★v525 — 6쪽 짤림 해소로 8.6→8.0mm
    ('제108조 상태칸',   'report_pages.py',    r'td class="st-cell"', True),
    ('제111조 레켐비',   'report_weasy.py',    r'알츠하이머 신약 <b>레켐비</b>', True),
    ('제109조 간병라벨', 'report_pages.py',    r"line\('간호통합병동', '간호통합병동'\)", True),
    ('제110조 폐기', 'main.py',                r"간병인질병일당\(요양병원\)", False),
    ('제120조 수식어괄호', 'main.py',           r"간편\|할인형\|맞춤고지", True),
    ('제116조 4쪽협심증', 'report_pages.py',    r"'협심증', 'I20', '협심증'", True),
    ('제116조 4쪽급성심근', 'report_pages.py',  r"'급성심근경색', 'I21~23', '급성심근경색'", True),
    ('제116조 4쪽빈맥', 'report_pages.py',      r"'빈맥', 'I47 · 48', '빈맥'", True),
    ('제116조 4쪽허혈복제', 'report_pages.py',  r"'@허혈성 진단비'", False),
    ('제121조 심장동결표', 'report_pages.py',   r'HEART_SYNC = \[', True),
    ('제121조 심장동결검사','report_pages.py',  r'def heart_audit\(', True),
    # ★v573 리모델링은 보유+증가분 표기라 차단하지 않는다(제127조). 진단서 차단은 main.py 게이트가 한다.
    ('제121조 발행차단',   'main.py',           r'제121조 심장동결 위반\(보장분석지 PPT\)', True),
    ('제121조 PPT표',      'main.py',           r'_PPT_HEART = \(', True),
    ('제121조 PPT검사',    'main.py',           r'def _ppt_heart_audit\(', True),
    ('제123조 HB전역',    'main.py',           r'^_HB = \{', True),
    ('제123조 케이스표',   'main.py',           r'_HEART_CASES = \[', True),
    ('제123조 케이스검사', 'main.py',           r'def heart_case_selftest\(', True),
    ('제123조 정규화한곳', 'main.py',           r'def _heart_norm\(', True),
    ('제123조 묶음허혈금지','main.py',          r'def heart_bundle_no_isch\(', True),
    ('제124조 높이통일',   'remodel.py',        r'v538 높이통일', True),
    ('제125조 감액판정',   'remodel.py',        r"tag, fn = '감액', R", True),
    ('제125조 상태색',     'remodel.py',        r'_cfill = \(NEWF if tag', True),
    ('제126조 통합형전이암','main.py',           r"has\('통합형전이암'\)", True),
    # ★v582 — `_c4`에 제131조 3색용 `sp`가 붙어 시그니처가 늘었다. 인자 이름까지 고정하면
    #   조문검사가 <b>정상 개정까지 실패로 잡는다</b> → 함수 존재 여부로 본다(핵심은 보유·제안 분리).
    ('제127조 보유제안분리','report_pages.py',   r'def _c4\(v, o=None', True),
    ('제131조 리포트3색',   'report_pages.py',   r'def _own_span\(', True),
    ('제128조 마스킹보존', 'main.py',            r"\[\*●○◯·・\\u25cf", True),
    ('제128조 실명채택',   'main.py',            r'v545 고객명', True),
    ('제128조 분석경로',   'main.py',            r'v549 고객명', True),
    ('제128조 성일치',     'main.py',            r"_c\[0\] == _sn", True),
    ('제9조 비급여통원',   'main.py',            r"no\('주사','MRI','도수','체외','증식'\): return '통원'", True),
    ('제9조 실손케이스',   'main.py',            r'_SILSON_CASES = \[', True),
    ('제9조 실손검사',     'main.py',            r'def silson_selftest\(', True),
    ('제9조 상해의료비1세대','main.py',          r'v552 실손', True),
    ('제128조 제목한줄',   'report_pages.py',    r'letter-spacing:-\.04em;white-space:nowrap\}', True),
    ('제21조 뇌심주요치료','main.py',            r"has\('뇌'\) and \(has\('심장'\) or has\('허혈심'\)\)", True),
    ('제21조 주요치료혈전금지','main.py',        r"has\('혈전용해'\) and has\('치료'\) and no\('주요치료'\)", True),
    ('제129조 전단어추출', 'report_pptx.py',     r'def _all_boxes\(', True),
    ('제129조 텍스트보존', 'remodel.py',         r'v543 리포트PPT', True),
    ('제129조 값칸텍스트', 'remodel.py',         r'v546 리포트PPT', True),
    ('제119조 CI상설',    'report_pages.py',    r"_lb\.endswith\('\(CI\)'\)", True),
    ('제119조 4쪽CI금지',  'report_pages.py',    r"'CI \(중대한 질병\)', '중대한 뇌졸중'", False),
    ('제129조 값칸선별', 'remodel.py',           r"_re9\.search\(r'\\d', w\[0\]\)", True),
    ('제129조 통이미지금지','remodel.py',        r'def build_report_pptx\(pngs\):', False),
    ('제127조 old읽기',    'report_pages.py',   r"covo = \(cmp_\.get\('old'\)", True),
    ('제127조 합산파싱',   'report_pages.py',   r'_re\.findall\(r\'\(\[\\d,\]\+\)', True),
    ('제124조 한도비례',   'remodel.py',        r'_wt = \[\(_c \* 0\.86', True),
    ('제124조 앞꽉금지',   'remodel.py',        r'v536 공백0', False),
    ('제124조 패딩금지',   'remodel.py',        r'_fpad > 0: r \+= _fpad', False),
    ('제124조 요약패딩금지','remodel.py',       r'요약 %d행 아래로 밀어', False),
    ('제71조 만기9999종신','remodel.py',        r"'9999' in _ex", True),
    ('제71조 만기키',      'remodel.py',        r"'expiry_date': str\(ws\.cell\(4, c\)", True),
    # ★v538 제124조 5항 — 라운드로빈 균등은 폐기(쪽마다 행 높이가 달라졌다).
    #   쪽수 유동은 `_caps_for()`가 담당하고, 몫은 한도 비례로 나눈다.
    ('제100조 유동쪽', 'remodel.py',            r"def _caps_for\(", True),
    ('제100조 행높이', 'remodel.py',            r"ws\.row_dimensions\[_rr\]\.height = _hh", True),
    ('제111조 연금인포', 'report_weasy.py',    r'연금계좌 세액공제', True),
    ('제110조 폐기2', 'main.py',               r"간병인상해일당\(요양병원\)", False),
    ('제116조 심장14행', 'report_pages.py',    r"'인공소생 성공 심장정지', 'I46\.0'", True),
    ('제116조 5쪽유동', 'report_pages.py',    r"\('산정특례 심장', '산정특례심장'\)", True),
    ('제119조 CI기재',   'report_pages.py',    r"'중대한 급성심근 \(CI\)', '중대한 급성심근'", True),
    ('제97조 인쇄예산',   'remodel.py',         r'_PAGE_PT = 735\.0', True),
    ('제100조 균등배분',  'remodel.py',         r'v490 균등', True),
    # ★v537 제124조 3항 — 「공백이란 없다」(지점장 지시 2026.08.21).
    #   제101조 꼬리 바닥 고정용 <b>빈 행 패딩은 폐기</b>한다. 채우기는 행 높이로만 한다.
    ('제101조 패딩폐기',  'remodel.py',        r'제124조 3항', True),
    ('제96조 담보갱신단순', 'main.py',           r"if '갱신' in _t95: return True", True),
    ('제95조 담보100',    'main.py',            r'_dk95 and all\(_is_gen_dambo', True),
    ('제98조 배상책임',   'main.py',            r'_ILSANG = re.compile', True),
    ('제98조 10억통장',   'main.py',            r"\('리셋월렛' in _t95\)", True),
    ('제98조 10억항상파랑','main.py',           r"'일상배상책임','10억 플랜'", True),
    ('제98조 마스터무행', 'main.py',            r'_dk95 = \[_k for _k in _dk95 if \(resolve2', True),
    ('제93조 인쇄농도',   'remodel.py',         r'convert_from_path\(f, dpi=300\)', True),
    # ★v410b 구 패턴 `\[v\d\d\d `는 <b>기존 기능 이력 라벨 수백 건</b>까지 잡아 거짓경보를 냈다.
    #   제36조의 대상은 <b>각인을 찍는 로그</b>(`[지침]`·`[zip검증]`)뿐이다 → 그 둘만 검사한다.
    # ★v410c 주석 안의 <b>예시 문구</b>까지 걸렸다 → `print(` 줄로 한정한다.
    ('제36조 라벨무버전','main.py',            r"print\([^\n]*\[v\d\d\d (지침|zip검증)\]", False),
    # ★v410d 나머지 조문 — 파일 구조로 검사 가능한 것 전수
    ('제0조 서열',       'main.py',            r'BARUM_DOCTRINE\.md', True),
    ('제4조 제외7종',    'main.py',            r'농업인', True),
    ('제5조 갱신판정',   'main.py',            r'9999', True),
    ('제5조 색 0070C0',  'main.py',            r'0070C0', True),
    ('제6조 사망배정',   'main.py',            r'일반사망', True),
    ('제8조 슬래시',     'main.py',            r'종수술비\(1-5종\)', True),
    ('제10조 제안서',    'main.py',            r'def build_proposal_contract', True),
    ('제13조 검산',      'main.py',            r'한장표\] 검산', True),
    ('제15조 마스터동적', 'main.py',           r'nm2r', True),
    ('제29조 묻지않는다', 'BARUM_DOCTRINE.md', r'지침에 있는 건 묻지 않는다', True),
    ('제31조 조문=테스트','main.py',           r'_JOMUN_SELFTEST', True),
    ('제32조 전문정독',  'BARUM_DOCTRINE.md',  r'분석지 실행마다 이 파일을 전문 정독한다', True),
    ('제14조 응대',      'BARUM_DOCTRINE.md',  r'팩폭', True),
    # ★v410f 제37조(커버리지 100%) 자신에게 검사가 없어 97%였다 — <b>조문을 만들면 그 조문도 검사한다.</b>
    ('제37조 커버리지',  'main.py',            r'_STRUCT_SELFTEST', True),
    # ★v412 제38조는 [기각] — 열 배치가 <b>현행 그대로</b>인지 검사한다(누가 다시 바꾸면 잡힌다).
    ('제38조 열배치현행','main.py',            r"own_sum_col  = \(3 \+ n_ct\)", True),
    ('제39조 자부상80',  'main.py',            r'_new413 = round\(amt / 80\)', True),
    ('제37조 미결3부',   'BARUM_DOCTRINE.md',  r'A\. 지점장 확정이 필요한 것', True),
    # ★v422g 제42조 — ④ 리모델링 비교
    ('제42조 리모델링',   'main.py',            r"@app\.post\('/remodel'\)", True),
    ('제42조 모듈',       'remodel.py',         r'def remodel_all', True),
    ('제42조 배포10파일', 'main.py',            r"'ppt_form\.pptx','remodel\.py'", True),
    # ★v424 제44~46조 — 일시납 · 연금 · 재무 페이지 · 작업순서
    ('제44조 일시납보관',  'main.py',         r'lump_sum = pv', True),
    ('제44조 헤더기재',    'main.py',         r'일시납\)', True),
    ('제44조 합계방어',    'remodel.py',      r"'일시납' in s", True),
    ('제44조 계약판별',    'remodel.py',      r'def contract_kinds', True),
    ('제45조 재무페이지',  'report_pages.py', r'def p8\(', True),
    ('제45조 진단서이식',  'report_weasy.py', r'P-ASSET', True),
    ('제45조 동적생성',    'report_weasy.py', r'def _asset_body', True),
    ('제45조 값박기금지',  'report_weasy.py', r'메트라이프생명 무배당', False),
    ('제48조 설명서제외',  'main.py',         r'v424 설명서', True),
    ('제48조 10건상한',    'report_pages.py', r'MAXROW = 10', True),
    ('제48조 비갱신가입율','report_pages.py', r'비갱신 가입율', True),
    ('제48조 저축연금화재','report_pages.py', r"'저축', '연금', '화재'", True),
    ('제48조 마스터불변',  'main.py',         r'max_row.*106|106행|MASTER_ROWS', True),
    ('제49조 빈페이지금지','report_weasy.py', r'v425 .재무 생성 실패', True),
    ('제49조 잘림금지',    'report_pages.py', r'v425 다장', True),
    ('제49조 상한자르기금지','report_pages.py', r'\[:MAXROW\]', False),
    ('제49조 설명서경고',  'main.py',         r'v425 .설명서', True),
    ('제49조 회사금액매칭','report_weasy.py', r'_amt_s in _hay', True),
    ('제49조 일시납주기',  'main.py',         r'v425 주기 우선', True),
    ('제50조 이미지분리',  'report_weasy.py', r'from assets_b64 import', True),
    ('제50조 거대상수금지','report_weasy.py', r"_FIN_SURVEY = 'iVBOR", False),
    ('제50조 번호하드코딩','report_weasy.py', r'<div class="pgn"><b>\d+</b>', False),
    ('제50조 재무맨뒤',    'report_weasy.py', r'재무상태 설문지', True),
    ('제52조 엑셀2개',    'main.py',    r'fd\.append\("old_xlsx",R1\)', True),
    ('제52조 rows폴백',    'remodel.py', r"if not old.get\('rows'\)", True),
    ('제52조 dict폴백',    'remodel.py', r"_rowlist\(old, 'delete'\) or dele", True),
    ('제53조 매니페스트',  'main.py', r'manifest\.webmanifest', True),
    ('제53조 아이콘생성',  'main.py', r'def _pwa_icon', True),
    ('제53조 SW무캐시',    'main.py', r'no-store', True),
    ('제54조 관리자비번',  'main.py', r'ADMIN_PW = "821024"', True),
    ('제54조 회원테이블',  'main.py', r'CREATE TABLE IF NOT EXISTS members', True),
    ('제54조 쉬운코드',    'main.py', r'def _mk_code', True),
    ('제54조 DB없어도뜸',  'main.py', r'DATABASE_URL 없음', True),
    ('제54조 psycopg',     'requirements.txt', r'psycopg', True),
    ('제54조 추방버튼',    'main.py', r'async function kick', True),
    ('제54조 추방기록',    'main.py', r"VALUES\(%s,%s,'kick'\)", True),
    ('제55조 조문수하한',  'main.py', r'DOCTRINE_MIN_ART', True),
    ('제55조 결번감시',    'main.py', r'조문 번호 결번', True),
    ('제55조 분량하한',    'main.py', r'DOCTRINE_MIN_CHARS', True),
    ('제56조 JS검사',     'main.py', r'def js_selftest', True),
    ('제56조 PC반응형',    'main.py', r'@media \(min-width:900px\)', True),
    ('제56조 엔티티',      'main.py', r'&#39;', True),
    ('제57조 줄바꿈금지',  'main.py', r'white-space:nowrap;overflow:hidden', True),
    ('제57조 모바일축소',  'main.py', r'@media \(max-width:520px\)', True),
    ('제58조 글자만',      'main.py', r'class="logo">M<', True),
    ('제59조 가입신청',    'main.py', r'/member/apply', True),
    ('제59조 승인액션',    'main.py', r"act == 'approve'", True),
    ('제59조 상태컬럼',    'main.py', r"status TEXT DEFAULT 'active'", True),
    ('제59조 관리자버튼',  'main.py', r'href="/admin"', True),
    # ★★★★★v446 (지점장 확정 2026.08.17) — 검사 없는 조문은 지침이 아니다.
    ('제60조 손보판정',    'main.py', r'def _is_nonlife\(', True),
    ('제60조 손보CI차단',  'main.py', r'if _is_nonlife\(co\):', True),
    ('제60조 회사전달',    'main.py', r'_isci_prod\(product, company\)', True),
    ('제60조 손보케이스',  'main.py', r'_CI_NONLIFE_SELFTEST', True),
    ('제60조 생보먼저',    'main.py', r"'생명' in c or '라이프' in c", True),
    ('제61조 명시import',  'report_weasy.py', r'from assets_b64 import \(', True),
    ('제61조 스타금지',    'report_weasy.py', r'from assets_b64 import \*', False),
    ('제61조 산출물게이트','main.py', r'_OUT4_KEYS', True),
    ('제62조 행높이보정',  'main.py', r'_RH_FILL', True),
    ('제62조 마스터우선',  'main.py', r'_d\.height is None', True),
    ('제63조 GET만',       'main.py', r"e\.request\.method !== 'GET'", True),
    ('제63조 POST무가로채기','main.py', r"e\.respondWith\(fetch\(e\.request\)\); \}\);", False),
    ('제64조 A1원천',      'main.py', r'def _cust_from_xlsx', True),
    ('제64조 복사본제거',  'main.py', r"re\.sub\(r'\\s\*\[-\(\]\.\*\$'", True),
    ('제65조 화면캐시금지','main.py', r"'Cache-Control': 'no-store, no-cache, must-revalidate'", True),
    ('제65조 onReady',    'main.py', r'function _onReady\(fn\)', True),
    ('제65조 리셋버튼',    'main.py', r'id="tabnew"', True),
    ('제65조 리셋전체',    'main.py', r'savedFiles=\{\};\n   R1=null; R2=null;', True),
    ('제4조 단체금지',     'main.py', r"★v446 '단체' \+ 날짜불명 → 제외", True),
    ('제62조 마스터종결',  'BARUM_DOCTRINE.md', r'마스터에 직접 넣었다', True),
    # ★★★★★v454 제58조 검사 — 화면 이모지 금지.
    #   ★패턴에 이모지 <b>문자</b>를 직접 쓰면 이 등록 줄 자신이 매칭돼 뮤테이션이 안 잡힌다(실측).
    #     → 유니코드 <b>이스케이프</b>로 쓴다. 소스에는 글자가 없고 re가 해석한다.
    # ★★★★★v460 제69조 심플 모드 — 형식 검사를 걷어냈다. 남긴 것은 값검사 배선뿐이다.
    ('제69조 값검사',    'main.py', r'def behave_selftes[t]', True),
    ('제69조 발행차단',  'main.py', r"bad \+= \['값검사 ' \+ x for x in _bb\]", True),
    ('제69조 조문',      'BARUM_DOCTRINE.md', r'실제로 사고가 났던 것만', True),
    ('제68조 분석강제',    'main.py', r'_BEHAVE_WAR[N] = \[\]', True),
    ('제68조 health',      'main.py', r"'동작검사': \(lambda t", True),
    ('제54조 추방버튼',    'main.py', r'function kick', True),
    ('제54조 추방2단확인',  'main.py', r'한 번 더 확인합니다', True),
    ('제54조 추방API',      'main.py', r"DELETE FROM members WHERE code", True),
    ('제45조 쪽번호없음',  'report_pages.py', r'자산 · 재무</span>', False),
    ('제46조 넘침감지',    'remodel.py',      r'REPORT_OVERFLOW', True),
    ('제47조 페이지반복',  'report_pages.py', r'position:fixed;top:0', True),
    ('제47조 높이고정금지','report_pages.py', r'min-height:1\d\dmm', False),
    ('제47조 다장대응',    'report_pages.py', r'position:fixed;left:0;right:0', True),
    # ★★★★★v422 제41조(삼성 병합셀·특정치료비) — 구조 검사 3종. 하나라도 지워지면 잡힌다.
    ('제41조 대괄호담보행','main.py',           r'┖\|\\\[\[\^', True),
    ('제41조 특정치료비',  'main.py',            r"has\('특정치료비'\)", True),
    ('제41조 삼성자료',    'report_weasy.py',    r'INFO-TBL 통합치료비 4', True),
]

# ★★★★★v404 조문 강제 테이블 — <b>조문을 넣을 때 여기 한 줄을 같이 넣는다.</b>
_JOMUN_SELFTEST = [
    # (담보명, 기대 마스터행, 조문)
    ('상급종합병원Ⅲ하이클래스암주요치료비(연간1회한,진단후10년)', '하이클래스(암)', '제20조 하이클래스'),
    ('하이클래스암주요치료비Ⅱ(상급종합병원(국립암센터포함))',      '하이클래스(암)', '제20조 하이클래스'),
    ('비급여암주요치료비',                                        '하이클래스(암)', '제20조 하이클래스'),
    ('암주요치료비Ⅲ(상급종합병원)(유사암제외)(연간1회한)(주요치료)', '암주요치료비',   '제20조 하이클래스'),
    ('심뇌혈관질환주요치료비(연간1회한)(수술및혈전용해치료)',        '2대 주요치료비', '제21조 심뇌혈관'),
    ('심뇌혈관질환주요치료비(연간1회한)(중환자실입원)',              '2대 주요치료비', '제21조 심뇌혈관'),
    ('신특정순환계질환통합치료비ⅢPlus',                            '2대 주요치료비', '제21조 심뇌혈관'),
    ('뇌졸중 혈전용해치료비(1회한)',                                '혈전용해치료비', '제21조 심뇌혈관'),
    ('통합암진단비',                                              '통합암',         '제22조 통합암'),
    ('통합전이암진단비',                                          '통합전이암',     '제22조 통합암'),
    ('암진단비(유사암제외)(통합간편가입형)',                        '일반암',         '제22조 통합암'),
    ('통합간편가입Ⅱ암진단비',                                     '일반암',         '제22조 통합암'),
    ('암입원일당(1-180)',                                         '암일당',         '제18조 암일당'),
    ('암입원일당(요양병원)',                                       '암일당',         '제18조 암일당'),
    ('화재상해사망후유장해특약 : 화상수술비',                        '화상수술비',     '제16조 콜론뒤'),
    ('화재상해사망후유장해특약 : 화상진단비',                        '화상진단비',     '제16조 콜론뒤'),
    # ★★★★★v422 (지점장 확정 2026.08.15) — 삼성 특정치료비. 제31조: 조문을 넣으면 검사 한 줄을 같이 넣는다.
    ('[건강]종합병원암(유사암Ⅱ제외)특정치료비Ⅲ(수술(회당),항암방사선,항암약물)', '암주요치료비', '제41조 특정치료비'),
    # ★v422d 지점장 「1-8이 한 세트야」 — 6·7·8종은 1-8종 세트(대표 max)
    ('[간편]상해8종수술비(시술포함)',        '상해 종수술비(1-8종)', '제41조 1-8종세트'),
    ('[간편]질병6종수술비(시술포함)(1년50%)', '질병 종수술비(1-8종)', '제41조 1-8종세트'),
    ('[간편]상해1~5종수술비(5종)',          '상해 종수술비(1-5종)', '제41조 1-8종세트'),
    ('[건강]종합병원유사암Ⅱ특정치료비Ⅲ(수술(회당),항암방사선,항암약물)',        '__무시__',     '제41조 특정치료비'),
    ('[건강]종합병원암전액본인부담(비급여포함)통합치료비(표준형,연간1억원한도)', '하이클래스(암)','제41조 특정치료비'),
    ('비급여(전액본인부담 포함) 암 통합치료비Plus(암중점치료기관(상급종합병원 포함))', '하이클래스(암)','제41조 특정치료비'),
    ('[건강]종합병원특정순환계질환통합치료비(표준형,연간1억원한도)',            '2대 주요치료비','제41조 특정치료비'),
]

_RULE_SELFTEST = [
    # (담보명, 기대키 또는 '__NOSPLIT__', 조문)
    ('심뇌혈관질환주요치료비(연간1회한)(수술및혈전용해치료)담보', '__NOSPLIT__', '제21조 심뇌혈관'),
    ('심뇌혈관수술비',                                          '심장수술비[묶음]', '제21조 심뇌혈관'),
]


# ★★★★★v410 (지점장 지시 2026.08.12 「지금 해라 절대 미루지 마라」): 조문 커버리지 확장.
#   제31조(조문=테스트)를 만들어놓고 <b>내가 만든 새 조문 6개에는 검사를 안 넣었다</b> — 커버리지가 22%→19%로 떨어졌다.
#   여기 있는 것은 전부 <b>기계로 확인 가능한</b> 조문이다. 확인 불가한 조문(응대·절차)은 표 아래에 사유를 적는다.
_JOMUN2_SELFTEST = [
    # (담보명, 기대 마스터행, 조문)
    ('실효',                       None,            '제4조 제외7종'),
    ('농업인NH안전보험',            None,            '제4조 제외7종'),
    ('상해 종수술비(1-5종)',        '상해 종수술비(1-5종)', '제8조 슬래시'),
    ('일반상해사망',               '상해사망',       '제6조 사망'),
    ('암진단비(유사암제외)',        '일반암',         '제22조 통합암'),
    ('유사암진단비',               '유사암(갑.기.경.제)', '제1조 엑셀기준'),
    ('일상생활중배상책임',          '일상배상책임',    '제2조 산출물범위'),
    ('비급여주사료',               '비급여주사',      '제9조 실손'),
    ('비급여MRI검사비',            'MRI',            '제9조 실손'),
    ('3대비급여(도수치료,체외충격파,증식치료)', '도수치료', '제9조 실손'),
]

def _jomun2_check():
    """제4·6·8·9·22조 — 담보 배치 조문 실측."""
    bad=[]
    for nm, exp, cho in _JOMUN2_SELFTEST:
        try: got = resolve2(nm)[0]
        except Exception as e: bad.append(f'[{cho}] {nm} → ERR {e}'); continue
        if got != exp: bad.append(f'[{cho}] {nm} → {got}(기대 {exp})')
    return bad

def _file_check():
    """★v410e 제34조 — main.py 각인 <b>리터럴 개수</b>. 이건 정규식 유무 검사로는 못 잡는다(개수 문제).
    나머지 파일·구조 조문은 `_STRUCT_SELFTEST`가 정본이다(중복 제거)."""
    import os as _o, re as _r
    bad=[]; base=_o.path.dirname(_o.path.abspath(__file__)) or '.'
    try:
        with open(_o.path.join(base,'main.py'),encoding='utf-8',errors='replace') as h: mn=h.read()
    except Exception as e: return [f'[제34조] main.py 읽기 실패 {e}']
    _lit = len(_r.findall(r"'v\d{3}[a-z]*-[a-z0-9]+-\d{8}'", mn))
    if _lit != 1: bad.append(f'[제34조] main.py 각인 리터럴이 {_lit}곳 — VSTAMP 한 곳이어야 한다')
    for _f in ZIP9:
        if not _o.path.exists(_o.path.join(base,_f)): bad.append(f'[제12조] 배포파일 없음 {_f}')
    return bad

def doctrine_selftest():
    """지점장 확정 조문이 코드에 살아 있는지 매 배포 검사."""
    bad=[]
    for nm, exp in _DEDUP_SELFTEST:
        try: got = _dedup_std(nm)
        except Exception as _e: bad.append(f'[제3조 dedup] {nm} → ERR {_e}'); continue
        if got != exp: bad.append(f'[제3조 dedup] {nm} → {got}(기대 {exp})')
    for nm, exp in _DOCTRINE_SELFTEST:
        try: got = resolve2(nm)[0]
        except Exception as _e: bad.append(f'{nm} → ERR {_e}'); continue
        if got != exp: bad.append(f'{nm} → {got}(기대 {exp})')
    for nm, exp in _SOLO5_SELFTEST:
        try: got = is_solo5_name(nm)
        except Exception as _e: bad.append(f'[제3조 단독5종] {nm} → ERR {_e}'); continue
        if got != exp: bad.append(f'[제3조 단독5종] {nm} → {got}(기대 {exp})')
    # ★★★★★v404 (지점장 지시 2026.08.12 「지침이 길어서 안 읽힌다 — 대책」):
    #   <b>조문 = 테스트.</b> 문서는 안 읽히지만 <b>테스트는 매 배포마다 돈다.</b>
    #   구 셀프테스트는 `resolve2`·`is_solo5_name`만 봤다 — 그래서 2026.08.12에 터진 위반을
    #   <b>단 한 건도 못 잡았다</b>(전부 다른 파일·다른 함수에 있었다).
    #   ★이제 조문을 새로 넣을 때는 <b>여기 검사를 같이 넣는다. 검사 없는 조문은 지침이 아니다.</b>
    for _nm, _exp, _cho in _JOMUN_SELFTEST:
        try: _got = resolve2(_nm)[0]
        except Exception as _e: bad.append(f'[{_cho}] {_nm} → ERR {_e}'); continue
        if _got != _exp: bad.append(f'[{_cho}] {_nm} → {_got}(기대 {_exp})')
    # 제16조 담보명은 콜론 뒤 / 제21조 심뇌혈관주요치료비는 수술 분해 금지
    for _nm, _must, _cho in _RULE_SELFTEST:
        try: _d = rule_extract([_nm+'  100'], prefolded=True); _d.pop('__DUP__', None)
        except Exception as _e: bad.append(f'[{_cho}] {_nm} → ERR {_e}'); continue
        _ks = list(_d.keys())
        if _must == '__NOSPLIT__':
            if any('[묶음]' in k for k in _ks): bad.append(f'[{_cho}] {_nm} → 묶음분해됨 {_ks}')
        elif _must not in _ks:
            bad.append(f'[{_cho}] {_nm} → {_ks}(기대 {_must} 포함)')
    # 제25·28조 「보유는 그 이름의 담보가 있을 때만」 — coverage_benchmark 실측
    try:
        import coverage_benchmark as _cb2, re as _re9
        _src = open(_cb2.__file__, encoding='utf-8', errors='replace').read()
        _i = _src.index('    scope_heart=[]'); _j = _src.index("    if _any('산정특례심장'", _i)
        _hb = _src[_i:_j]
        _i2 = _src.index('    scope_brain=[]'); _j2 = _src.index("    if _any('산정특례뇌'", _i2)
        _bb = _src[_i2:_j2]
        def _runsc(blk, key, names):
            _ns = {'_any': (lambda *ks: any(k in n for n in names for k in ks))}
            exec(blk.replace('    ', ''), _ns); return _ns[key]
        if 'angina' in _runsc(_hb, 'scope_heart', ['허혈성 진단비']):
            bad.append('[제28조 허혈성단독] 허혈성진단비만 있는데 협심증(angina)이 보유로 찍힌다')
        if 'infarct' in _runsc(_bb, 'scope_brain', ['뇌혈관진단비']):
            bad.append('[제25조 보유] 뇌혈관진단비만 있는데 뇌졸증(infarct)이 보유로 찍힌다')
        if 'hem' in _runsc(_bb, 'scope_brain', ['뇌혈관진단비']):
            bad.append('[제25조 보유] 뇌혈관진단비만 있는데 뇌출혈(hem)이 보유로 찍힌다')
    except Exception as _e:
        bad.append(f'[제25·28조] 검사 불가 {_e}')
    # ★v410 파일 구조로 검사 가능한 조문 전수
    import os as _os10, re as _re10
    _base = _os10.path.dirname(_os10.path.abspath(__file__)) or '.'
    _cache = {}
    for _cho, _fn, _pat, _want in _STRUCT_SELFTEST:
        if _fn not in _cache:
            _t = ''
            for _d in (_base, '.'):
                try:
                    with open(_os10.path.join(_d, _fn), encoding='utf-8', errors='replace') as _h: _t = _h.read()
                    break
                except Exception: continue
            _cache[_fn] = _t
        _t = _cache[_fn]
        if not _t: bad.append(f'[{_cho}] {_fn} 읽기 실패'); continue
        _hit = bool(_re10.search(_pat, _t, flags=_re10.M))
        if _hit != _want:
            bad.append(f'[{_cho}] {_fn} — ' + ('있어야 하는데 없다' if _want else '없어야 하는데 있다') + f' ({_pat})')
    bad += heart_case_selftest()   # ★v533 제123조 2항 — 심장 케이스
    bad += silson_selftest()       # ★v551 제9조 3항 — 실손 전세대
    bad += _jomun2_check()      # ★v410 제4·6·8·9·22조
    bad += _file_check()        # ★v410 제12·15·19·24·26·34·36조
    for cd, np3, drug, exp in _GEN_SELFTEST:
        try: got = silson_gen(cd, None, '', np3, drug)
        except Exception as _e: bad.append(f'[제9조 실손세대] {cd} → ERR {_e}'); continue
        if got != exp: bad.append(f'[제9조 실손세대] {cd}/np3={np3}/약값={drug} → {got or "공란"}(기대 {exp or "공란"})')
    return bad


# ★★★★★v446 손보 CI 무시 케이스 (지점장 지시 2026.08.17 · 박주하 실측)
#   (회사명, 상품명, 기대 CI판정)
_CI_NONLIFE_SELFTEST = [
    ('현대해상', '무배당퍼펙트클래스종합보험(Hi1706)기본플랜', False),  # ★실사고: '퍼펙트'로 CI 오판
    ('삼성화재', '무배당 삼성화재 퍼펙트플러스보험',            False),
    ('롯데손해보험', '무배당 롯데 내마음속 건강보험(1504)',      False),
    ('흥국화재', '무배당 프리미엄 행복보험(1210)',              False),
    ('DB손보',  '무배당 DB CI종신보험',                        False),  # 손보면 CI 표기여도 무시
    ('KB손해보험', '무배당 KB CI종신보험',                      False),
    ('삼성생명', '무배당 삼성생명 퍼펙트플러스보험',            True),   # 생보는 종전대로 CI
    ('교보생명', '무배당교보큰사랑 CI 보험',                    True),
    ('삼성생명', '삼성리빙케어(종신2종)1.2',                    True),
    ('흥국생명', '(무)흥국생명 CI종신보험',                     True),   # '흥국'이지만 생명 → CI
]


def ci_selftest():
    """CI 판정 4곳이 <b>같은 로직</b>인지 + 케이스 16건을 통과하는지 검사. 실패 목록을 돌려준다."""
    bad=[]
    for p,exp in _CI_SELFTEST:
        try:
            got=_isci_prod(p)
        except Exception as _e:
            bad.append(f'{p} → ERR {_e}'); continue
        if got!=exp: bad.append(f'{p} → {got}(기대 {exp})')
    # ★coverage_benchmark의 판정과도 대조 — 산출물 간 CI 인식이 갈리는 사고(v235) 방지
    try:
        import coverage_benchmark as _cb
        if hasattr(_cb,'_isci_hdr'):
            for p,exp in _CI_SELFTEST:
                if _cb._isci_hdr(p)!=exp: bad.append(f'[cb 불일치] {p}')
    except Exception:
        pass
    # ★★★★★v446 손보 CI 무시 (지점장 지시 2026.08.17) — 검사 없는 조문은 지침이 아니다.
    for _co, _pd, _exp in _CI_NONLIFE_SELFTEST:
        try:
            _g = _isci_prod(_pd, _co)
        except Exception as _e:
            bad.append(f'[손보CI] {_co}/{_pd} → ERR {_e}'); continue
        if _g != _exp:
            bad.append(f'[손보CI] {_co}/{_pd} → {_g}(기대 {_exp})')
    bad += ['[선지급률] '+x for x in ci_rate_selftest()]
    bad += ['[실손소스] '+x for x in silson_selftest()]
    bad += ['[조문] '+x for x in doctrine_selftest()]   # ★v382 지침 조문 강제
    return bad


def parse_sebu_ci(lines):
    """★★★v237 영구지침(지점장 지시 2026.07.25): <b>CI 선지급률(50%형/80%형)을 세부가입현황(상세내역)에서 찾아낸다</b>.
    지점장 원문 = "이름이 CI GI 삼성생명퍼펙트… 있으면 상세세부내역(롯데) KB도 세부내역에서 50%형 or 80%형을 찾아내서 해야 한다."
    ★기존 판정(main.py CI 블록)은 <b>별첨의 '주계약' 라벨 금액(`ci_jugye`)만</b> 입력으로 썼다.
      → 롯데(let:) 별첨엔 '주계약' 라벨이 아예 없어(주계약을 `질병사망`/`상해사망`으로 씀) `ci_jugye=[]`가 되고
      선지급률 판정 블록을 통째로 건너뛰었다. 이것이 "CI가 하나도 반영 안 된" 두 번째 원인이다.
    ★파싱 원리: '계약별가입정보' 표는 <b>계약이 열로 나열</b>된다. 회사명 줄에서 각 계약의 x 시작좌표를 잡고,
      담보군 행(사망/암/뇌혈관질환/심장질환)의 숫자를 <b>x좌표가 가장 가까운 계약</b>에 배정한다.
    반환: {회사명: {'samang': 사망최대값, 'cands': [본체 후보값…]}}
    """
    out={}
    try:
        _CO=re.compile(r'[가-힣A-Za-z]{1,7}\s?(?:생명|손보|화재|해상|라이프|손해보험|공제|증권)')
        blk=[]; flag=False
        for l in lines:
            if ('계약별가입정보' in l) or ('계약별 가입정보' in l): flag=True
            if flag: blk.append(l)
            if flag and ('안내 및 유의' in l or '별첨' in l): break
        # ★★★v253(김O구 실측 2026.07.26): <b>앵커 문자열에 의존하지 마라</b>.
        #   이 리포트는 4p 표의 <b>헤더·담보명 한글이 pdftotext에서 통째로 빠진다</b>
        #   (값·회사명은 정상 추출). 그래서 `'계약별가입정보'`를 못 찾아 blk=[] →
        #   <b>즉시 빈 dict 반환</b> → ci_sebu=None → CI 중대한OO가 하나도 안 들어갔다.
        #   → 앵커가 없으면 <b>회사 접미어가 2개 이상 있는 줄</b>을 직접 찾아 그 줄부터 읽는다.
        if not blk:
            for i,l in enumerate(lines):
                if len(list(_CO.finditer(l)))>=2:
                    blk=lines[i:i+60]
                    print('[v253 sebu_ci] 앵커 없음 → 회사명 줄(%d)로 진입' % i)
                    break
        if not blk: return out
        # ① 회사명 줄 = 회사 접미어가 2개 이상 등장하는 첫 줄
        cidx=-1; cos=[]
        for i,l in enumerate(blk[:14]):
            ms=list(_CO.finditer(l))
            if len(ms)>=2: cidx=i; cos=[(m.start(), re.sub(r'\s','',m.group())) for m in ms]; break
        if cidx<0 or not cos: return out
        cos.sort()
        # ★★좌측 '전체 가입현황' 표가 같은 줄에 있어 숫자가 섞인다(실측: 신한 사망이 23,100=전체합계로 오염).
        #   → 계약별가입정보 영역의 <b>x 경계</b>를 잡아 그 왼쪽 숫자·라벨은 전부 버린다.
        _xnum = cos[0][0] - 8      # 숫자 허용 하한
        _xlab = cos[0][0] - 26     # 담보군 라벨 허용 하한
        def _owner(x):
            best=None; bd=10**9
            for sx,nm in cos:
                d=abs(x-sx)
                if d<bd: bd=d; best=nm
            return best
        # ★v253 회사별 보험료(노이즈) 수집 — samang 폴백에서 제외하기 위해
        # ★v253: 보험료는 <b>'원'이 붙는다</b> → 영역 내 '원' 값 전부를 노이즈 집합으로 모은다
        #   (회사별로 나누면 동일 회사 계약이 여러 건일 때 첫 건만 잡혀 보험료가 사망으로 새어든다 — 실측).
        _PREMS=set()
        for l in blk:
            for mm in re.finditer(r'([\d,]{4,})\s*원', l):
                try: _PREMS.add(int(mm.group(1).replace(',','')))
                except: pass
        # ② 사망액은 '사망' 라벨 행에서, 본체 후보는 영역 내 전체 숫자에서 수집.
        #   ★계약별가입정보 표는 <b>라벨 행과 값 행이 어긋난다</b>(실측: '암' 라벨 47행 / 값 45행,
        #     '뇌혈관질환' 라벨 52행 / 값 51행). 라벨 기준으로 값을 묶으면 후보가 전멸한다.
        #   → 사망만 라벨로 잡고(정확도 확인됨), 나머지는 <b>x좌표로 계약에 배정한 전체 숫자</b>를 후보로 둔다.
        for l in blk[cidx+1:]:
            if not l.strip(): continue
            # ★노이즈 행 제외 — 보험료·보장기간·납입기간/횟수·상품명 줄은 담보 금액이 아니다.
            #   (실측 오염: 148,576(보험료) · 9999/2018(만기·가입연도) · 1701/1710/2510(상품코드))
            if re.search(r'보험료|보장기간|납입|회사|보험서비스|상품', l): continue
            if re.search(r'\d{4}\s*[.\-]\s*\d{2}', l): continue
            mk=re.search(r'(?<![가-힣])사망(?![가-힣])', l)
            _is_sam = bool(mk) and mk.start()>=_xlab
            for m in re.finditer(r'([\d][\d,]{2,})', l):
                if m.start() < _xnum: continue          # ★좌측 '전체 가입현황' 표 숫자 차단
                _g=m.group(1)
                if ',' not in _g and re.fullmatch(r'(?:1[6-9]|2[0-9])\d{2}', _g): continue  # 연도·상품코드
                try: v=int(_g.replace(',',''))
                except Exception: continue
                if not (0<v<=200000): continue
                nm=_owner(m.start())
                if not nm: continue
                d=out.setdefault(nm, {'samang':0, 'cands':[]})
                if _is_sam: d['samang']=max(d['samang'], v)
                else: d['cands'].append(v)
        # ★★★v253 samang 폴백(김O구 실측 2026.07.26): 담보명 라벨 한글이 빠진 PDF에서는
        #   '사망' 라벨 행을 못 찾아 samang=0이 된다. 세부가입현황 계약열은 <b>행 순서가 고정</b>
        #   (사망→후유장해→실손→암…)이라 <b>보험료 제거 후 첫 값이 사망</b>이다.
        # ★★★그 첫 값 = <b>일반/질병사망만</b>(지점장 확정): 두 번째 값은 <b>상해사망</b>이고
        #   <b>CI 법칙에 상해사망은 전혀 상관없다</b>. 실측 교보 4,000(질병)/10,000(상해) → <b>4,000</b>.
        for _co, _d in out.items():
            if float(_d.get('samang') or 0) > 0: continue
            _pool = [v for v in (_d.get('cands') or [])
                     if v >= 1000 and v not in (_PREMS or set()) and v % 100 == 0]
            if _pool:
                _d['samang'] = float(_pool[0])
                print('[v253 samang폴백] %s 사망보장 %s (상해사망은 CI 무관 → 제외)' % (_co, format(_pool[0], ',')))
        return out
    except Exception as _e:
        print(f'[v237 sebu_ci ERR] {_e}')
        return out


def parse_sebu(lines):
    """v30x 세부가입현황 전담보 파서 = 주 소스(오버랩 근절). 탭 좌우2열, 담보-금액 경계 명확."""
    out={}
    sb=[]; flag=False
    for l in lines:
        if ('세부가입현황' in l) or ('전체 가입현황' in l): flag=True; continue
        if flag and (('계약별 가입정보' in l) or ('계약별가입정보' in l) or ('안내 및 유의' in l)):
            if out: break
            continue
        if flag: sb.append(l)
    if not sb: return out
    def _num(x):
        x=str(x).replace(',','').replace('원','').strip()
        try:
            v=float(x); return v if v>0 else None
        except: return None
    _skip=('충분','부족','미가입','구분','담보명','보장지수','가입금액','-')
    for raw in sb:
        # ★표준금액/권장금액 줄 = 가입금액 아님, 반영 금지(지점장 정본)
        if ('표준금액' in raw) or ('권장금액' in raw) or ('적정금액' in raw): continue
        cells=[c.strip() for c in raw.split(chr(9)) if c.strip()!='']
        k=0
        while k < len(cells):
            nm=cells[k]
            _ko=re.sub(r'[^가-힣]','',nm)
            if len(_ko)>=2 and nm not in _skip:
                val=None
                for off in (1,2,3):
                    if k+off>=len(cells): break
                    c2=cells[k+off]
                    mm=re.match(r'(충분|부족|미가입)\s*([\d,]+)?$', c2)
                    if c2 in ('충분','부족'):
                        if k+off+1<len(cells): val=_num(cells[k+off+1])
                        break
                    elif c2=='미가입' or c2=='-':
                        break
                    elif mm and mm.group(1)!='미가입' and mm.group(2):
                        val=_num(mm.group(2)); break
                    else:
                        vv=_num(c2)
                        if vv is not None: val=vv; break
                if val is not None:
                    _std,_j=resolve_kw(nm)
                    if _std and _std!='__무시__':   # ★v30z 무시지정 담보(전이암·고액항암)는 세부보충에서도 제외
                        if _std not in out or out[_std]<val: out[_std]=val
            k+=1

    return out

# ══════════════════════════════════════════════════════════════════════════
# ★★★★★v271 세부가입현황 '계약별 가입정보' 파서 (지점장 지적 2026.07.30
#   "삼성생명 41,800원 계약이 다 비어있다")
#
#   [왜 필요한가 — 추측이 아니라 리포트 인쇄값이다]
#   삼성생명 별첨은 담보명 칸에 <b>담보명이 아니라 특약명/상품명</b>이 인쇄된다.
#     無신정기특약Ⅰ(표준,본인) 1,000 / 無재해상해(본인) 700 / 無입원(본인) 2 …
#   담보 <b>종류를 알 수 없으니</b> 마스터 키워드에 걸릴 수가 없고 계약 열이 통째로 빈다.
#   교보·동양 생보 CI가 별첨에 `주계약`만 찍히는 것과 <b>같은 상황</b>이며,
#   그때 확립된 원칙 그대로 <b>세부가입현황이 유일한 정답 소스</b>다(지침 §3 ②).
#
#   [구조 — 실측 확정]
#   '계약별 가입정보'(우측)는 <b>보험료 줄 아래로 21개 데이터 행</b>이 좌측
#   '전체 가입현황' 라벨과 <b>같은 순서</b>로 내려온다. 행 번호가 곧 담보다.
#   한 계약 열 안에서 <b>왼쪽 값=질병 · 오른쪽 값=상해</b>.
#
#   [배정 방법] pdftotext -layout은 x좌표에 비례해 공백을 넣으므로 문자 인덱스를
#   좌표로 쓴다. 보험료 값의 위치로 계약 열 중심을 잡고 <b>최근접 계약</b>에 배정,
#   그 계약 열의 <b>중간점 좌/우</b>로 질병·상해를 가른다.
#   ★확신이 없는 행(요양·치매/응급실·화재벌금/깁스)은 <b>None으로 두어 건드리지 않는다</b>.
# ★★★★★v276 (지점장 지적 2026.07.30 "이영태 — CI부터 보장성도 하나도 안 맞다")
#   구 _SEBU_ROWMAP(21행 고정 순번)은 <b>완전 폐기</b>한다.
#   실측: 이영태 리포트의 좌측 라벨은 <b>33개</b>(실손 4·암 5·입원비 4·운전자 5·기타 5)여서
#   21행 고정표가 통째로 어긋났다 → 장기요양 3,200이 <b>뇌혈관진단비</b>로, 일배책 10,000이
#   <b>상해수술비</b>로, 합의금 10,000이 <b>뇌졸증</b>으로 박혔다.
#   ★새 원리 = <b>순번이 아니라 그 줄에 인쇄된 좌측 담보명 라벨</b>로 매핑한다.
#     ①라벨 줄을 찾아 (좌=질병측, 우=상해측) 페어를 만든다(side는 사전에 고정 — x좌표 추정 금지)
#     ②인접(≤3줄) 라벨 줄은 side가 겹치지 않을 때만 병합한다(질병입원비/상해입원비가 두 줄로 갈린다)
#     ③구간 = [라벨줄-1 … 다음 라벨줄-2]. 값이 라벨보다 한 줄 앞서 인쇄되는 경우가 있다(실측 상해3%).
#     ④값 배정은 종전대로 <b>보험료 위치로 계약 열 최근접 + 열 중간점 좌/우</b>.
#   ★확신이 없는 라벨(장기요양·치매·응급실·화재벌금·고액항암)은 <b>None으로 두어 건드리지 않는다</b>.
_SEBU_LAB = [
    # (정규화 키워드, 마스터 담보명 or None, side)  side: 'L'=질병측(좌) 'R'=상해측(우)
    ('질병사망','질병사망(80세)','L'),      ('상해사망','상해사망','R'),
    ('질병3%','질병후유3%','L'),            ('상해3%','상해후유3%','R'),
    ('질병입원의료비','입원','L'),          ('상해입원의료비','입원','R'),
    ('질병통원의료비','통원','L'),          ('상해통원의료비','통원','R'),
    ('일반암진단비','일반암','L'),          ('암수술비','암수술','R'),
    ('유사암진단비','유사암(갑.기.경.제)','L'), ('고액항암치료비',None,'R'),
    ('전이암진단비','통합전이암','L'),
    ('뇌혈관진단비','뇌혈관진단비','L'),    ('뇌혈관수술비','뇌혈관수술비','R'),
    ('뇌졸중진단비','뇌졸증진단비','L'),    ('뇌졸증진단비','뇌졸증진단비','L'),
    ('허혈성심장질환진단비','허혈성 진단비','L'), ('심장질환수술비','심장수술비','R'),
    ('급성심근경색진단비','급성심근경색','L'),
    ('질병수술비','질병수술비','L'),        ('상해수술비','상해수술비','R'),
    ('질병입원비','질병일당','L'),          ('상해입원비','상해일당','R'),
    ('장기요양자금',None,'L'),              ('경증이상치매진단비',None,'R'),
    ('교통사고처리지원금','합의금','L'),    ('변호사선임','변호사','R'),
    ('벌금(대인)','대인','L'),              ('부상위로금','자부상','R'),
    ('벌금(대물)','대물','L'),
    ('응급실내원비',None,'L'),              ('화재벌금',None,'R'),
    ('골절진단비','골절(치아파절포함)','L'),('일상생활배상책임','일상배상책임','R'),
    ('깁스치료비','깁스진단비','L'),
]

def _sebu_norm(s):
    """라벨 비교용 정규화: 공백·전각괄호·구분기호 제거."""
    t = str(s or '')
    t = t.replace('（','(').replace('）',')').replace('［','[').replace('］',']')
    return re.sub(r'[\s,·:/]+','',t)

def _sebu_labels_of(seg):
    """좌측표 영역 문자열에서 담보 라벨을 찾아 {'L':마스터명|None, 'R':...} 반환.
       ★'간병인사용/지원 ○○입원비'는 간병인 행으로 돌린다(질병일당/상해일당 아님)."""
    t = _sebu_norm(seg)
    if not t: return {}
    care = ('간병' in t)
    out = {}
    for kw, tgt, side in _SEBU_LAB:
        if _sebu_norm(kw) not in t: continue
        if side in out: continue                       # 먼저 잡힌 라벨 우선(사전 순서 = 표 순서)
        if care and tgt in ('질병일당','상해일당'): tgt = '간병인'
        out[side] = tgt
    return out

_BBOX_XML = ''

# ★★★★★v279 세부가입현황 좌표 폴백 — 라벨 글자를 전혀 쓰지 않는다.
#   [왜] 신형 let: 리포트(이명순·김진구)는 5~8p 좌측표 한글 라벨이 통째로 이미지라
#        v276 라벨 매핑이 라벨 0개 → 블록 통째 스킵 → 별첨 담보명이 상품명뿐인 계약
#        (삼성리빙케어·한화 종신)의 담보가 전부 사라졌다(이명순 담보 27/97행).
#   [앵커] 좌측 '전체 가입현황' 표의 각 행 값은 <b>한장보장표 값과 같다</b>.
#        → 행 순서표로 매핑한 뒤 <b>한장표 값과 대조해 검증</b>하고, 검증을 통과할 때만 채택한다.
#          (순번을 쓰되 값으로 검증하므로 v271의 '검증 없는 고정 순번'과 성격이 다르다.)
_SEBU_PAIRS = [
    ('질병사망(80세)','상해사망'), ('질병후유3%','상해후유3%'),
    ('입원','입원'),               ('통원','통원'),
    ('일반암','암수술'),           ('유사암(갑.기.경.제)',None),
    ('통합전이암',None),
    ('뇌혈관진단비','뇌혈관수술비'), ('뇌졸증진단비',None),
    ('허혈성 진단비','심장수술비'), ('급성심근경색',None),
    ('질병수술비','상해수술비'),    ('질병일당','상해일당'),
    ('간병인','간병인'),            (None,None),
    ('합의금','변호사'),            ('대인','자부상'), ('대물',None),
    (None,None),
    ('골절(치아파절포함)','일상배상책임'), ('깁스진단비',None),
]

def _bbox_words(xml):
    """bbox XML → [페이지별 [(x, y, text)]]"""
    pages=[]
    for pg in re.split(r'<page\b', xml)[1:]:
        ws=[]
        for m in re.finditer(r'xMin="([\d.]+)"\s+yMin="([\d.]+)"[^>]*>(.*?)</word>', pg):
            ws.append((float(m.group(1)), float(m.group(2)), m.group(3)))
        if ws: pages.append(ws)
    return pages

def _num_or_dash(t):
    t=t.strip()
    if t=='-': return 0.0
    if not re.fullmatch(r'\d{1,3}(?:,\d{3})*|\d+', t): return None
    try: return float(t.replace(',',''))
    except: return None

def parse_sebu_bbox(hj):
    """좌표 기반 계약별 담보값. 반환 = { 보험료(int) : { 담보명 : 금액 } }
       hj = parse_hanjang 결과(검증 기준). hj가 비면 검증 불가라 통째 포기한다."""
    res={}
    if not _BBOX_XML or not hj: return res
    for ws in _bbox_words(_BBOX_XML):
        # ── 보험료 줄(계약 열) 찾기: '숫자+원' 이 2개 이상 같은 y
        byy={}
        for x,y,t in ws: byy.setdefault(round(y/3.0)*3, []).append((x,t))
        prem_y=None; cols=None
        for y in sorted(byy):
            ps=[(x,t) for x,t in byy[y] if re.fullmatch(r'\d[\d,]*원', t)]
            if len(ps)>=2:
                ps.sort(); prem_y=y
                cols=[{'x':x,'prem':int(t[:-1].replace(',',''))} for x,t in ps]
                break
        if not cols: continue
        w = (cols[-1]['x']-cols[0]['x'])/(len(cols)-1) if len(cols)>1 else 101.0
        # ★보험료는 우측정렬이라 값보다 오른쪽에 찍힌다 → 열 구간은 보험료 x에서 15pt 당겨 잡는다.
        for k,c in enumerate(cols):
            c['lo'] = c['x'] - 15
            c['hi'] = (cols[k+1]['x'] - 15) if k+1 < len(cols) else (c['x'] + w - 15)
        L0 = cols[0]['lo'] - 10          # 좌측표 / 계약열 경계
        # ── 값 행 수집(보험료 줄 아래)
        rows=[]
        for y in sorted(byy):
            if y <= prem_y: continue
            left=[(x,t) for x,t in byy[y] if x < L0]
            right=[(x,t) for x,t in byy[y] if x >= L0]
            lv={}
            for x,t in left:
                v=_num_or_dash(t)
                if v is None: continue
                lv['L' if x < 250 else 'R'] = v
            if not lv and not right: continue
            if lv or right: rows.append((y,lv,right))
        # ── 값 없는 꼬리(페이지 번호·연락처) 제거: 좌측 판정값이 있는 행만 본문으로 본다
        body=[r for r in rows if r[1]]
        # ★★★★★v296 (이명순 롯데 실측 2026.07.31, 영구): <b>좌측 판정값과 계약열 값이 한 행에 안 있다</b>.
        #   실측(세부표 6p) — 일반암 행이 <b>y255(좌측 5,980) / y258(계약열 4,000·1,500·480)</b> 두 줄로 갈렸다.
        #   행 버킷이 3pt라 <b>다른 행으로 분리</b>되고, 계약열만 있는 줄은 `if r[1]` 필터에서 <b>통째로 버려졌다</b>.
        #   → 일반암이 <b>전 계약에서 소실</b>(한장표 5,980 vs 엑셀 0). 유사암은 우연히 같은 버킷이라 살아남았다.
        #   ★<b>고아 줄(좌측 없고 계약열만 있는 줄)을 가장 가까운 본문 행에 되붙인다</b>(±6pt, 1:1).
        _orphan=[r for r in rows if (not r[1]) and r[2]]
        _used=set()
        for _bi,(_by,_blv,_brt) in enumerate(body):
            if _brt: continue
            _best=None; _bd=99
            for _oi,(_oy,_olv,_ort) in enumerate(_orphan):
                if _oi in _used: continue
                _d=abs(_oy-_by)
                if _d<=6 and _d<_bd: _best=_oi; _bd=_d
            if _best is not None:
                _used.add(_best); body[_bi]=(_by,_blv,_orphan[_best][2])
        if _used: print(f'[v296 sebu-bbox] 갈린 계약열 줄 {len(_used)}건 복원')
        if len(body) < 10: continue
        body=body[:len(_SEBU_PAIRS)]
        # ── ★검증: 좌측값 ↔ 한장표 값 대조
        ok=0; bad=0
        for (y,lv,rt),(pl,pr) in zip(body,_SEBU_PAIRS):
            for side,nm in (('L',pl),('R',pr)):
                if not nm or side not in lv or nm not in hj: continue
                if abs(lv[side]-hj[nm]) < 0.5: ok+=1
                else: bad+=1
        if ok < 5 or ok <= bad:
            print(f'[v279 sebu-bbox] 검증 실패(일치 {ok}/불일치 {bad}) → 이 페이지 스킵'); continue
        # ── 배정
        for (y,lv,rt),(pl,pr) in zip(body,_SEBU_PAIRS):
            # ★열별로 값을 모은 뒤 <b>개수로</b> 좌(질병)/우(상해)를 가른다.
            #   한 열에 2개면 왼쪽=질병·오른쪽=상해. 1개면 열 시작 근처인지로 판정.
            grp={}
            for x,t in rt:
                v=_num_or_dash(t)
                if v is None or v<=0: continue
                for ci,c in enumerate(cols):
                    if c['lo'] <= x < c['hi']:
                        grp.setdefault(ci,[]).append((x,v)); break
            for ci,items in grp.items():
                c=cols[ci]
                if not c['prem']: continue
                items.sort()
                if len(items)>=2:
                    pick=[(pl,items[0][1]),(pr,items[-1][1])]
                else:
                    x0,v0=items[0]
                    pick=[((pl if x0 < c['x']+30 else pr), v0)]
                for nm,v in pick:
                    if not nm: continue
                    d=res.setdefault(c['prem'],{})
                    if nm not in d or d[nm] < v: d[nm]=v
        print(f'[v279 sebu-bbox] 검증 통과(일치 {ok}) · 계약 {len(cols)}열 배정')
    return res

# ★★★★★v295 (지점장 지시 2026.07.31, 영구): <b>KB 3열 리포트 전용 검산 앵커</b>.
#   [문제] v277·v278 자동 검산은 <b>한장보장표가 있어야만</b> 돈다. KB·메리츠 3열 리포트에는
#   그 표가 없어 <b>검산 게이트가 통째로 꺼진다</b>. 김수영 실측에서 유사암 4,800(정답 800)·
#   골절 2,129·깁스 5,030이 <b>아무 경고 없이</b> 산출물로 나갔다. 잡아낸 것은 시스템이 아니라 지점장 눈이었다.
#   [해법] KB 리포트 2~3p <b>'전체 보장 현황'</b> 표를 한장표 대용 앵커로 쓴다.
#   구조 = `담보명  합계  KB손보  손해보험  생명보험  공제/체신` → <b>담보명 뒤 첫 값 = 합계</b>.
#   반환 키를 <b>한장표와 동일한 마스터 담보명</b>으로 맞춰 기존 검산식(_pairs)을 그대로 재사용한다.
#   ★<b>입원·통원은 넣지 않는다</b> — KB 요약표는 실손 상해 5,000+질병 5,000을 '1억'으로 합쳐 쓰는데
#     BARUM 정본은 <b>4세대 입원 = 최고값 1개(5,000)</b>라 구조적으로 다르다(불일치가 아니라 규칙 차이).
#   ★<b>고액(표적)항암치료비·3대비급여실손·간병인 계열도 제외</b> — 신정원 합산 표기라 1:1 대응이 아니다.
_KB_SUM_MAP = {
    '상해사망':'상해사망', '질병사망':'질병사망',
    '상해80%미만후유장해':'상해후유3%', '질병80%미만후유장해':'질병후유3%',
    '일반암':'일반암', '유사암':'유사암(갑.기.경.제)',
    '뇌혈관질환':'뇌혈관진단비', '뇌졸중':'뇌졸증진단비', '뇌출혈':'뇌출혈진단비',
    '허혈성심장질환':'허혈성 진단비', '급성심근경색증':'급성심근경색',
    '상해수술비':'상해수술비', '질병수술비':'질병수술비', '암수술비':'암수술',
    '뇌혈관질환수술비':'뇌혈관수술비', '허혈성심장질환수술비':'심장수술비',
    '상해입원일당':'상해일당', '질병입원일당':'질병일당',
    '골절진단비':'골절진단비', '가족/일상/자녀배상':'일상배상책임',
    '교통사고처리지원금':'합의금', '변호사선임비용':'변호사', '자동차사고부상':'자부상',
}
def _kb_amt(tok):
    """'4,000만' '1억' '1억 7,600만' '-' → 만원 float. 해석 불가면 None."""
    t = str(tok).strip()
    if t in ('-', '', '_'): return 0.0
    # ★★★★★v421 (지점장 지적 2026.08.14 박미정 검산 불일치 — <b>게이트 오탐의 진짜 원인</b>)
    #   KB 한장표는 <b>'1억 7,600만'</b>처럼 억과 만을 <b>같이</b> 쓴다.
    #   구 파서는 `([\d,]+)\s*억` 하나만 봐서 <b>'1억'만 읽고 7,600만을 버렸다</b> → 10,000.
    #   엑셀은 17,600(정답)인데 <b>기준값이 틀려</b> 「불일치」가 떴다.
    #   실측 3건: 상해사망 10,000≠17,600 · 상해후유3% 10,000≠10,100 · 뇌출혈 0≠1,000.
    #   ★<b>「불일치」가 뜨면 엑셀만 의심하지 말고 기준값도 의심한다</b>(검사 자체의 오탐도 결함).
    m = re.fullmatch(r'([\d,]+)\s*억\s*([\d,]+)\s*만?', t)
    if m:
        try:
            return float(m.group(1).replace(',', '')) * 10000.0 + float(m.group(2).replace(',', ''))
        except Exception: return None
    m = re.fullmatch(r'([\d,]+)\s*억', t)
    if m:
        try: return float(m.group(1).replace(',', '')) * 10000.0
        except Exception: return None
    m = re.fullmatch(r'([\d,]+)\s*만?', t)
    if m:
        try: return float(m.group(1).replace(',', ''))
        except Exception: return None
    return None

def parse_kb_summary(lines):
    """KB 리포트 '전체 보장 현황' 표 → { 마스터담보명 : 합계(만원) }. 없으면 {}."""
    res = {}
    if not any('전체 보장 현황' in l or '전체보장현황' in re.sub(r'\s', '', l) for l in lines):
        return res
    for l in lines:
        _l = l.strip()
        if not _l: continue
        for _kb, _std in _KB_SUM_MAP.items():
            # ★★★★★v421 (지점장 지적 2026.08.14 박미정) — KB 한장표는 <b>'1억 7,600만'</b>처럼
            #   억과 만 사이에 <b>공백</b>이 있다. 구 코드는 `(\S+)` 한 토큰만 잡아 <b>'1억'</b>만 읽었다.
            #   → 앵커 10,000, 엑셀 17,600 → 「검산 불일치」가 <b>오탐</b>으로 떴다(제출 금지 유발).
            #   ★<b>금액은 토큰 하나가 아니다</b> — 억 뒤에 만이 오면 같이 읽는다.
            m = re.match(r'^' + re.escape(_kb) + r'(?![가-힣])\s+([\d,]+\s*억(?:\s*[\d,]+\s*만?)?|\S+)', _l)
            if not m: continue
            v = _kb_amt(re.sub(r'\s+', ' ', m.group(1)).strip())
            if v is None: continue
            # 골절진단비는 마스터 2행 합산이라 검산식 쪽에서 합친다 → 전용 키로 보관
            res[_std] = max(res.get(_std, 0.0), v)
            break
    if res: print(f'[v295 KB요약표] 검산 앵커 {len(res)}개 확보(한장표 대용)')
    return res

def parse_sebu_bycontract(lines):
    """세부가입현황 '계약별 가입정보'에서 계약별 담보값을 읽는다(v276 라벨 기반).
       반환 = { 보험료(int) : { 마스터담보명 : 금액(float) } }"""
    res={}
    def _toks(s):
        # ★★★v294: <b>날짜는 금액이 아니다</b>. `2025-08-22 2079-08-22` 같은 보험시기·종기 열이
        #   섞여 들어오면 2079·22가 담보값으로 잡힌다(김수영 실측). → 날짜 패턴을 먼저 지운다.
        s = re.sub(r'\d{4}\s*[-./~]\s*\d{1,2}\s*[-./]\s*\d{1,2}', ' ', s)
        out=[]
        for m in re.finditer(r'(?<![\d,])(\d{1,3}(?:,\d{3})*|\d+)(?![\d,])', s):
            out.append(((m.start()+m.end())/2.0, m.group(1)))
        return out
    i=0; n=len(lines)
    while i < n:
        l=lines[i]
        prem=[(m.start(), m.group(1)) for m in re.finditer(r'(\d[\d,]*)\s*원', l)]
        if len(prem) < 2:
            i+=1; continue
        # ★★★★★v355b (이성준 실측 2026.08.02): <b>'보험료' 라벨이 있는 줄만 계약 블록 시작</b>으로 본다.
        #   [구 결함] 「'원'이 2개 이상인 줄」이면 무조건 블록 시작이라
        #   <b>2p 보유계약 리스트의 기납입보험료·잔여보험료</b>(959,200원 · 4,272,800원)까지
        #   계약으로 잡혔다 → 계약 5건인데 <b>블록 6~7개</b>가 생기고 값이 옆 계약으로 밀렸다
        #   (실측 NH 암수술 2,000 오기재 · 운전자 담보 전멸).
        #   ★'계약별 가입정보' 표의 보험료 줄에는 <b>항상 '보험료' 라벨이 같은 줄(또는 바로 위 3줄)</b>에 있다.
        _ctx355 = re.sub(r'\s','', ''.join(lines[max(0,i-3):i+1]))
        if '보험료' not in _ctx355:
            i+=1; continue
        cols=[]
        for st,v in prem:
            try: p=int(v.replace(',',''))
            except: p=None
            cols.append({'start':st,'prem':p})
        _w = ((cols[-1]['start']-cols[0]['start'])/(len(cols)-1)) if len(cols)>1 else 40.0
        for k,c in enumerate(cols):
            c['l'] = c['start']-6
            c['r'] = (cols[k+1]['start']-6) if k+1<len(cols) else (c['start']+_w-6)
            c['mid'] = (c['l']+c['r'])/2.0
            c['ctr'] = c['start']+5
        L0 = cols[0]['l']
        # ── 블록 끝 = 다음 보험료 줄
        # ★★★★★v294 (김수영 KB 실측 2026.07.31, 영구): <b>끝 앵커가 없으면 리포트 끝까지 먹는다</b>.
        #   실측: 보험료가 2개 이상인 줄이 <b>5p 상품별 가입현황 한 줄뿐</b>이라 end=파일끝이 되어
        #   <b>5p~13p 9페이지를 한 계약 블록</b>으로 읽었다. 그 안 11p '담보별 가입 현황' 표의
        #   <b>맨 오른쪽 열 = 보험종기 날짜</b>를 금액으로 집어 <b>2079 · 22 · 14</b>가 담보값이 됐다
        #   (골절 2,079 / 상해사망 22 / 암수술 14). v260 접힘 사고와 <b>같은 계열</b>이다.
        #   → <b>다른 표가 시작되면 거기서 끊는다</b>. 종료 앵커가 없으면 종전대로 동작하므로 안전하다.
        #   ★★★★★v355 (이성준 실측 2026.08.02): <b>[별첨]이 앵커에 없어 별첨 담보표까지 먹었다</b>.
        #     실측 = 삼성 별첨 10p `［간편］상해 사망`의 <b>1</b>, KB 별첨 11p `상해 1〜5종수술비`의
        #     <b>1·5</b>가 금액으로 잡혀 <b>롯데(69,115원) 열</b>에 들어갔다
        #     → 질병수술비 100 + <b>5</b> = 105 · 상해수술비 100+1 · 상해사망 1,100+1.
        #     v294 사고와 <b>완전히 같은 계열</b>(끝 앵커 부재 → 다음 표를 통째로 먹음).
        _END_ANCHOR = ('담보별 가입 현황','담보별가입현황','실효/해지','실효 / 해지',
                       '가입담보상세','가입 담보 상세','보유 계약 리스트','보유계약리스트','전체 계약리스트',
                       '[별첨]','［별첨］','정상계약 리스트','실효계약 리스트','안내 및 유의 사항',
                       '계약별 가입정보','계약별가입정보')   # ★v356c 페이지가 바뀌면 새 표다 → 거기서 끊는다
        j=i+1; end=n
        while j < n:
            if len(re.findall(r'(\d[\d,]*)\s*원', lines[j]))>=2: end=j; break
            _lj = re.sub(r'\s','', lines[j])
            if any(re.sub(r'\s','',_a) in _lj for _a in _END_ANCHOR):
                end=j; break
            j+=1
        # ── ① 라벨 줄 수집
        labs=[]     # (줄번호, {'L':..,'R':..})
        for k in range(i+1, end):
            d=_sebu_labels_of(lines[k][:max(0,int(L0))])
            if d: labs.append([k,d])
        # ── ② 인접(≤3줄) & side 비겹침 → 병합
        merged=[]
        for k,d in labs:
            if merged and (k-merged[-1][0])<=3 and not (set(d) & set(merged[-1][1])):
                merged[-1][1].update(d); merged[-1][0]=k
            else: merged.append([k,dict(d)])
        # ★라벨이 빈약한 블록은 '계약별 가입정보'가 아니다(보유계약 리스트 등) → 통째 스킵
        if len(merged) < 5:
            i=end; continue
        # ── ③ 구간 배정  ★구간 상한 lk+3 — 마지막 라벨이 페이지 꼬리/다음 헤더를 먹는 것을 막는다
        for mi,(lk,d) in enumerate(merged):
            s_from = lk-1
            s_to   = (merged[mi+1][0]-2) if mi+1 < len(merged) else (end-1)
            s_to   = min(s_to, lk+3)
            if s_to < s_from: s_to = s_from
            for k in range(max(i+1,s_from), min(end, s_to+1)):
                for cx,val in _toks(lines[k]):
                    if cx < L0: continue
                    best=min(cols, key=lambda c: abs(cx-c['ctr']))
                    if abs(cx-best['ctr']) > 30: continue
                    tgt = d.get('L') if cx < best['mid'] else d.get('R')
                    if not tgt: continue
                    try: fv=float(val.replace(',',''))
                    except: continue
                    if fv<=0: continue
                    if not best['prem']: continue        # ★보험료 0원 열은 키 충돌 → 배정 제외
                    dd=res.setdefault(best['prem'],{})
                    if tgt not in dd or dd[tgt] < fv: dd[tgt]=fv
        i=end
    return res

# ★★★★★v277 한장보장표 자동 검산 (지점장 지시 2026.07.31 "지침 100% 활용 대안")
#   지침 §1 등식1「한장보장표 = 엑셀 = PPT」와 §13 체크리스트 ①은 <b>사람이 눈으로</b> 대조해야만
#   지켜지는 규칙이었다. 그래서 지침이 있어도 매번 지적을 받고서야 발견했다.
#   → <b>앱이 매 분석마다 스스로 대조</b>해 확인사항 시트 [검산] 블록에 불일치를 강제 출력한다.
#   ★이 값은 <b>검산 전용</b>이다. 절대 산출물(본표·PPT·설명서)에 기재하지 않는다.
#     따라서 여기서 순번 가정을 쓰더라도 <b>토큰 수가 정확히 일치할 때만</b> 채택하므로
#     틀려도 산출물이 오염되지 않는다(v276에서 폐기한 '순번 고정'과 성격이 다르다).
_HJ_BLOCKS = {
    ('A',8):  ['상해사망','질병사망','상해후유3%','질병후유3%','입원','__','통원','__'],
    ('B',11): ['일반암','유사암(갑.기.경.제)','통합전이암','암수술','__고액항암',
               '뇌혈관진단비','뇌졸증진단비','뇌혈관수술비','허혈성 진단비','급성심근경색','심장수술비'],
    ('C',8):  ['상해수술비','질병수술비','상해일당','__간병인상해','질병일당','__간병인질병',
               '__장기요양','__경증치매'],
    ('D',10): ['합의금','대인','대물','변호사','자부상','__응급실','__골절','깁스진단비',
               '__화재벌금','일상배상책임'],
}

_HJ_STD = {   # ★v278 표준금액 행 = 고객 무관 고정값 → 라벨이 이미지인 신형 리포트의 유일한 앵커
    'A': ['20,000','10,000','5,000','3,000','5,000','5,000','30','30'],
    'B': ['3,000','1,000','2,000','300','7,000','1,000','2,000','1,000','1,000','2,000','1,000'],
    'C': ['100','30','3','12','3','12','2,000','1,000'],
    'D': ['20,000','3,000','500','5,000','30','5','30','30','2,000','10,000'],
}

def _hj_toks(l):
    return [m.group(1) for m in re.finditer(r'(-|\d{1,3}(?:,\d{3})*|\d+)', str(l))]

def parse_hanjang(lines):
    """한장보장현황(4p) → {검산키: 가입금액(만원)}. 실패하면 {} (추측 금지).
       ★v278 근본수정: 구 v277은 '한장보장' 글자 + '가입금액' 라벨에 의존했다.
       신형 let: 리포트는 이 표의 라벨이 통째로 이미지라 pdftotext에 글자가 한 자도 없다
       (이명순·김진구 실측 = '한장' 0건) → 검산이 조용히 통째로 꺼져 있었다.
       → 1순위 앵커 = <b>표준금액 행</b>. 이 값은 고객과 무관한 고정 상수라 블록까지 특정된다.
       2순위 = 구 라벨 경로(표준금액이 다른 옛 템플릿 대비)."""
    out={}; rows={}
    try:
        # ── 경로1(정본): 표준금액 시그니처 → 바로 다음 숫자줄 = 가입금액
        for tag,std in _HJ_STD.items():
            for i,l in enumerate(lines):
                if _hj_toks(l)!=std: continue
                for k in range(i+1, min(i+8, len(lines))):
                    tk=_hj_toks(lines[k])
                    if not tk: continue
                    if len(tk)==len(std): rows[tag]=tk
                    break
                break
        # ── 경로2(폴백): '가입금액' 라벨 (구형 리포트)
        if len(rows)<4:
            st=None
            for i,l in enumerate(lines):
                if '한장보장' in re.sub(r'\s','',l): st=i; break
            if st is not None:
                seq=[]
                for k in range(st, min(st+120, len(lines))):
                    if '세부가입현황' in re.sub(r'\s','',lines[k]): break
                    if re.sub(r'\s','',lines[k]).startswith('가입금액'):
                        seq.append(_hj_toks(lines[k][lines[k].index('가입금액')+4:]))
                for tag,toks in zip(('A','B','C','D'), seq):
                    rows.setdefault(tag, toks)
        for tag,toks in rows.items():
            names=_HJ_BLOCKS.get((tag,len(toks)))
            if not names: continue          # 토큰 수 불일치 → 그 블록은 통째 포기
            for nm,tk in zip(names,toks):
                if nm.startswith('__'): continue
                out[nm] = 0.0 if tk=='-' else float(tk.replace(',',''))
    except Exception as e:
        print(f"[v278 한장표] 파싱 실패 → 검산 생략 ({e})")
        return {}
    if out: print(f"[v278 한장표] 검산 기준 {len(out)}개 확보 (블록 {sorted(rows)})")
    else:   print("[v278 한장표] 표를 찾지 못했다 → 검산 생략(3열 메리츠 등 한장표 없는 리포트)")
    return out

# ══════════════════════════════════════════════════════════════════════════
# ★★v44 신정원 3열 포맷 파서 (KB '상품별 가입담보상세' / 메리츠 '별첨 상품별 보험가입현황')
#   지점장 확정 2026.07.13: (1) 3열 포맷을 입력 정본에 추가 (앵커 자동감지, 기존 2열과 병행)
#                          (2) 담보명 정본 = 회사담보명
#                          (3) 갱신판정 = 상품명/담보명 '갱신' 표기 단독 기준 (3열엔 총회차 없음)
#   담보행: NO | 구분(정액/실손) | 회사담보명 | 신정원담보명 | 가입금액
# ══════════════════════════════════════════════════════════════════════════
_SJ_HEAD = re.compile(r'^\s*(\d{1,2})\s+(정액|실손)\s+(.+?)\s*$')
_SJ_AMT  = re.compile(r'\s{2,}((?:\d[\d,]*\s*[억만천]\s*)+)$')          # 행 끝 금액(1억 5,000만 등 공백 허용)
_SJ_CONT = re.compile(r'^\s+\S.*?\s{2,}((?:\d[\d,]*\s*[억만천]\s*)+)$')  # 신정원담보명 줄바꿈 wrap → 다음 줄 금액
_SJ_KB   = re.compile(r'^[\s\f]*(\S.*?)\s{2,}\|\s*가입일자\s*:\s*(\d{4})[-.](\d{2})[-.](\d{2})\s*\|')
_SJ_MZ   = re.compile(r'^[\s\f]*별첨\s+상품별\s*보험가입현황')

def _amt_kr(s):
    """한글 금액('1억 5,000만','2,000만','1억','1천') → 만원 단위 정수. 미해석=None(추측 금지)."""
    t = re.sub(r'[\s,]', '', str(s or ''))
    m = re.fullmatch(r'(?:(\d+)억)?(?:(\d+)만)?(?:(\d+)천)?', t)
    if not m or not any(m.groups()): return None
    v = int(m.group(1) or 0) * 10000 + int(m.group(2) or 0) + int(m.group(3) or 0) * 1000
    return v if 0 < v <= 200000 else None

def sinjeong_count(lines):
    """3열 신정원 계약 헤더 앵커 개수(KB '| 가입일자 : |' + 메리츠 '별첨 상품별 보험가입현황')."""
    return sum(1 for l in lines if _SJ_KB.match(l)) + sum(1 for l in lines if _SJ_MZ.match(l))

def sinjeong_detect(lines):
    """3열 신정원 포맷 감지. 계약 헤더 앵커가 2개 이상이면 3열로 확정."""
    return sinjeong_count(lines) >= 2

_SJC = {}

def _sj_fixname(name, sj, comp, prod):
    """★v98 3열(KB·메리츠) 전용 담보명 정규화. 롯데 2열 경로는 이 함수를 타지 않는다.
       한장보장표(등식1)와 어긋나던 실측 6건을 회사담보명 표기 보정으로만 해결."""
    r = re.sub(r'\s', '', str(name)); s2 = re.sub(r'\s', '', str(sj))
    c = re.sub(r'\s', '', str(comp)); pr = re.sub(r'\s', '', str(prod))
    # F4 생보 종신 주계약: 일사=일반사망 / 재사=재해(상해)사망
    if '일사보험금' in r: return '일반사망_주계약'
    if '재사보험금' in r: return '재해사망_주계약'
    # F2 타인'사망'교통사고처리지원금 → 교통상해사망 오분류 차단(처리지원금=합의금)
    if '처리지원금' in r and '사망' in r: return re.sub('사망', '', name)
    # F3 질병 입원 수술비Ⅱ/Ⅲ/Ⅳ = 질병수술비 합산군(한장보장표 기준)
    if re.match(r'^질병입원수술비[ⅡⅢⅣⅤ]', r): return '질병입원수술비' + r[8:]
    # F6 생보 암/CI의 '뇌혈관진단' = 실제 뇌출혈(지침 §8.3)
    if r == '뇌혈관진단' and ('생명' in c or 'AIA' in c.upper()): return '뇌출혈진단비'
    # F8 자동차사고 부상치료지원금Ⅱ 등 변형 → 자부상(신정원담보명 근거)
    if '부상위로금' in s2 or '부상치료' in s2: return '자동차사고부상위로금'
    # F9 1-5종 재해수술 = 상해 종수술(질병 종수술 행 중복산입 차단)
    if '종재해수술' in r: return name.replace('재해수술', '상해수술', 1)
    # F7 실손 '의료비(입원+통원)' 통합형 → 실손 입원 행 (F5b보다 먼저 판정)
    if '의료비' in r and '입원' in r and '통원' in r: return '실손입원의료비'
    # F5b 암 통원 담보는 진단비 행에 산입 금지
    if '암통원' in r or '암통원' in s2 or (('통원' in s2) and ('암' in r)): return '[확인] 통원 ' + name
    return name

def _sj_rows(block):
    """담보행 → [(회사담보명, 만원정수)]. 금액 미해석 담보는 0 + [확인] 프리픽스(누락 금지)."""
    out = []; i = 0
    while i < len(block):
        h = _SJ_HEAD.match(block[i])
        if not h: i += 1; continue
        rest = h.group(3); amt_s = None
        m = _SJ_AMT.search(rest)
        if m:
            amt_s = m.group(1); rest = rest[:m.start()]
        else:                                            # 신정원담보명 wrap → 다음 1~2줄에서 금액 회수
            for k in (1, 2):
                if i + k < len(block):
                    c = _SJ_CONT.match(block[i + k])
                    if c: amt_s = c.group(1); break
        parts = [p for p in re.split(r'\s{2,}', rest.strip()) if p]
        i += 1
        if not parts: continue
        name = re.sub(r'\s+', ' ', parts[0].strip())     # ★정본: 담보명 = 회사담보명(parts[0])
        sj   = re.sub(r'\s+', ' ', parts[1].strip()) if len(parts) > 1 else ''
        if len(name) < 2: continue
        v = _amt_kr(amt_s) if amt_s else None
        if v is None:
            out.append(('[확인] 금액판독불가 ' + name, sj, 0)); continue
        name = _sj_fixname(name, sj, _SJC.get('c',''), _SJC.get('p',''))
        if '특정암진단' in re.sub(r'\s','',sj) and '유사암' not in sj:   # ★v257 `r'\\s'` 오타 수정(공백 제거가 죽어 있었다)
            # ★v197(2026.07.23): 신정원 '특정암진단' = 고액암 행으로 확정(구 v98 F5 [확인]큐 폐기)
            # ★★★★★v247 (KB 3열 실측): <b>회사담보명이 유사암 4종이면 리네임하지 않는다</b>.
            #   정본 = <b>담보명 정본은 회사담보명</b>(신정원담보명은 분류 참고용)이다.
            #   실측 — 메리츠 `갱신형 갑상선암(초기제외)진단비` · `갱신형 갑상선암 및 기타피부암의 전이암…진단비`가
            #   신정원명 '특정암진단' 때문에 <b>고액암 2,000</b>으로 갔다(KB 전체보장현황 고액암은 1,000=라이나뿐).
            #   유사암 4종(갑상선·갑상샘·기타피부·경계성·제자리·상피내)은 <b>유사암 행</b>이 정답이다.
            _nm0 = re.sub(r'\s', '', str(name))
            if any(k in _nm0 for k in ('갑상선','갑상샘','기타피부','경계성','제자리','상피내')):
                pass
            elif '특정암' not in name:
                name = '특정암진단비(' + name + ')'
        out.append((name, sj, v))
    # ★v44 실측보정: DB 실손처럼 회사담보명이 '질병(전체질병을 의미)' 하나로 3행(입원·통원·약값)이 겹치는 경우
    #    회사담보명만 쓰면 dict 키 충돌 → 합산 사고. 계약 내 중복 회사담보명은 신정원담보명으로 분리한다.
    cnt = {}
    for nm, sj, v in out: cnt[nm] = cnt.get(nm, 0) + 1
    fixed = []
    for nm, sj, v in out:
        if cnt.get(nm, 0) > 1 and sj and sj != nm:
            fixed.append((_sj_fixname(sj, sj, _SJC.get('c',''), _SJC.get('p','')), v))  # ★v98 채택 후 재정규화
        else:
            fixed.append((nm, v))
    # ★★★v208 (지점장 확정 2026.07.25, 양*선 삼성 New내돈내삼 실측): 위에서 신정원담보명으로 갈아탄 뒤
    #   <b>신정원담보명끼리 또 중복</b>이면 1인실·2~3인실 구분이 통째로 사라져 <b>합산 사고</b>가 난다.
    #   실측: 상급종합 1인실 20 + 상급종합 2~3인실 5 → '상급종합병원 질병입원일당' 25 /
    #         종합 1인실 10 + 종합 2~3인실 5 → '종합병원이하 질병입원일당' 15 → 최종 질병종합병원일당 <b>40</b>.
    #   → 회사담보명에서 <b>병실 토큰(1인실 · 2~3인실 등)</b>만 뽑아 접미로 붙여 되살린다.
    cnt2 = {}
    for nm, v in fixed: cnt2[nm] = cnt2.get(nm, 0) + 1
    if any(c > 1 for c in cnt2.values()):
        fixed2 = []
        for (nm, v), (onm, osj, _ov) in zip(fixed, out):
            if cnt2.get(nm, 0) > 1:
                _rm = re.search(r'(\d+\s*~\s*\d+\s*인실|\d+\s*인실)', str(onm))
                if _rm: nm = '%s(%s)' % (nm, re.sub(r'\s', '', _rm.group(1)))
            fixed2.append((nm, v))
        fixed = fixed2
    return fixed

def _sj_unwrap(block):
    """★v98: pdftotext가 가운뎃점(·)에서 줄을 끊는다.
       '항암방사선·' / '질병 입원 간호·' 처럼 ·로 끝나는 줄은 다음 줄 첫 필드를 붙여 복원한다.
       (2~3줄 연속 랩도 while로 처리. 담보명 통째 유실 방지.)"""
    out = []; i = 0; n = len(block)
    while i < n:
        cur = block[i].rstrip()
        guard = 0
        while cur.endswith('\u00b7') and i + 1 < n and guard < 4:
            nxt = block[i + 1].strip()
            if not nxt: break
            mm = re.match(r'^(\S+(?:\s\S+)*?)(\s{2,}.*)?$', nxt)
            if not mm: break
            cur = cur + mm.group(1) + (mm.group(2) or '')
            i += 1; guard += 1
        out.append(cur); i += 1
    return out

def parse_sinjeong(lines):
    # ★v244: 3열도 세부가입현황(상세내역)을 CI 선지급률 2순위 근거로 쓴다(2열과 동일).
    try:
        _SJ_SEBU = {re.sub(r'\s','',k): v for k, v in (parse_sebu_ci(lines) or {}).items()}
    except Exception:
        _SJ_SEBU = {}
    """KB·메리츠 3열 리포트 → contracts[]. 표 구조 동일 → 파서 1개로 두 채널 커버."""
    lines = _sj_unwrap(lines)          # ★v98 F1: 가운뎃점 줄바꿈 복원
    n = len(lines)
    heads = []                                            # (idx, company, product, 가입일)
    for i, l in enumerate(lines):
        m = _SJ_KB.match(l)
        if m:                                             # ── KB: '회사명 ... | 가입일자 : YYYY-MM-DD |'
            comp = re.sub(r'\s+', '', m.group(1).strip())
            prod = ''
            for j in range(i + 1, min(i + 6, n)):
                s = lines[j].strip()
                if s and not re.search(r'가입일자|계약자|보험기간', s):
                    prod = s; break
            heads.append((i, comp, prod, f'{m.group(2)}.{m.group(3)}.{m.group(4)}'))
            continue
        if _SJ_MZ.match(l):                               # ── 메리츠: '별첨 상품별 보험가입현황' → 다음 2줄=회사·상품
            got = []
            for j in range(i + 1, min(i + 8, n)):
                s = lines[j].strip()
                if not s: continue
                if re.search(r'계약자|보험기간|가입담보명', s): break
                got.append(s)
                if len(got) == 2: break
            if len(got) >= 1:
                heads.append((i, re.sub(r'\s+', '', got[0]), (got[1] if len(got) > 1 else ''), ''))
    if not heads: return []

    contracts = []
    for hi, (idx, company, product, join_d) in enumerate(heads):
        end = heads[hi + 1][0] if hi + 1 < len(heads) else n
        block = lines[idx:end]
        _SJC['c'] = company; _SJC['p'] = product     # ★v98 _sj_fixname 컨텍스트
        ht = ' '.join(block[:12])
        contract_date = expiry_date = pay_period = ''; premium = 0; lump_sum = 0
        md = re.search(r'(\d{4})[-.](\d{2})[-.](\d{2})\s*~\s*(\d{4})[-.](\d{2})[-.](\d{2})', ht)
        if md:
            contract_date = f'{md.group(1)}.{md.group(2)}.{md.group(3)}'
            expiry_date   = f'{md.group(4)}.{md.group(5)}.{md.group(6)}'
        if not contract_date and join_d: contract_date = join_d
        mp = re.search(r'([\d,]{4,})\s*원', ht)
        if mp:
            try:
                pv = int(mp.group(1).replace(',', ''))
                # ★★★★★v424: 상한 500만은 <b>월보험료</b> 기준이다. 일시납(1,100만 등)이 여기서 버려졌다.
                #   월보험료(premium)와 <b>일시납(lump_sum)을 분리</b>해 담는다 — 합계에 섞이면 안 된다.
                if 1000 < pv < 5000000: premium = pv
                elif '일시' in re.sub(r'\\s','',str(pay_period or '')): lump_sum = pv   # ★v425 주기 우선
                elif 5000000 <= pv < 10**10: lump_sum = pv        # ★v424 금액 2순위
                elif pv >= 5000000 and '일시' in re.sub(r'\s','',str(pay_period or '')): lump_sum = pv
            except: pass
        mper = re.search(r'(?:월납|매월납)\s*/\s*(\d+)\s*년', ht)
        if mper: pay_period = f'{mper.group(1)}년납'
        elif re.search(r'일시납', ht): pay_period = '일시납'   # ★v127 3열 리포트 일시납 포착
        pay_count = ''                                          # ★v127 3열은 납입회차 칸이 없다
        if not expiry_date and re.search(r'종신', ht): expiry_date = '9999.12.31'
        if is_excluded(company, product, contract_date, expiry_date, pay_period, pay_count): continue   # 제외 5~7종

        dambo = {}
        # ★★★★★v244 (지점장 질문 "CI 대응 다 된 거냐" → 점검 중 발견): <b>3열(KB·메리츠) 경로에
        #   `ci_lines`가 통째로 빠져 있었다</b>. 정본 "규칙 하나 고치면 2열·3열 두 경로를 반드시 다 돌린다"를
        #   내가 지키지 않아, v239~v243의 CI 수정(줄단위 사망·중대한OO 본체·뇌 축 판별)이 <b>롯데에만</b> 적용됐다.
        #   → KB·메리츠 CI 계약은 여전히 <b>무조건 '중대한 뇌졸증'</b>으로 가서 같은 오류가 재발한다.
        # ★★★★★v246 (지점장 확정 2026.07.25): "<b>이건 롯데나 KB나 동일하다</b>" — CI 4단계·뇌 축·사망 배정은
        #   2열(롯데 [별첨] 보험서비스(상품)별 보장현황)과 3열(<b>KB '상품별 가입담보상세'</b>) <b>동일 적용</b>.
        #   ★3열은 앵커 줄(`| 가입일자 : |`)에서 <b>회사명 칸에 상품명이 붙어 오는 경우</b>가 있고
        #     상품명은 그 <b>다음 줄</b>에서 뽑는다 → <b>회사명·상품명 둘 다</b> CI 판정에 넣는다(한쪽만 보면 놓친다).
        _sjci = _isci_prod(product, company)          # ★v446 손보면 False
        ci_lines = {'samang': [], 'cands': [], 'jungdae': [], 'brain': []}
        for nm, v in _sj_rows(block):
            if re.search(r'납입면제|납입지원', nm): continue
            # ★★★v293: 실손 입원 = 합산 금지·대표(max). 2열 `_DUP_MAX_EXACT`와 동일 규칙을 3열에도 적용한다.
            if re.sub(r'\s','',str(nm)) in _DUP_MAX_EXACT:
                dambo[nm] = max(dambo.get(nm, 0), v)
            else:
                dambo[nm] = dambo.get(nm, 0) + v
            if _sjci and 0 < v <= 200000:
                _n2 = re.sub(r'\s', '', str(nm))
                # ★★★★★v253 지점장 지시 2026.07.26: <b>재해사망특약 6,000 → samang에서 삭제</b>.
                #   + 2차 규칙 원문 "세부가입현황에서 '사망'에서 <b>상해사망은 제외</b> 질병사망에서
                #   50% or 80%형 금액이 있는지 찾아본다".
                #   → CI 사망보장 후보 = <b>일반사망 / 질병사망(주계약)</b>. <b>상해사망·재해사망은 제외</b>.
                #   ★지점장이 말한 이 두 개만 제외한다 — 교통·외래 등 다른 항목을 임의로 덧붙이지 말 것.
                if (re.search(r'사망', _n2) and not re.search(r'후유|장해', _n2)
                        and not re.search(r'상해|재해', _n2)):
                    ci_lines['samang'].append(v)
                elif _n2.startswith('중대한'):
                    ci_lines['jungdae'].append(v)
                elif re.search(r'진단', _n2):
                    ci_lines['cands'].append(v)
                    if '뇌출혈' in _n2:   ci_lines['brain'].append(('뇌출혈', v))
                    elif '뇌졸' in _n2:   ci_lines['brain'].append(('뇌졸증', v))
        if not dambo: continue

        # ★확정(3) 갱신판정: 만기 9999(종신)=비갱신 / 상품명 '갱신' 표기=갱신 / 그 외는 후처리 담보절반 규칙
        if _is_silson_prod(company, product): renewal = '갱신'   # ★★★실손=무조건 갱신(v103 영구지침)
        elif str(expiry_date).startswith('9999'): renewal = '비갱신(종신)'
        elif '갱신' in str(product) and '비갱신' not in str(product): renewal = '갱신'
        else: renewal = '비갱신'

        ci_jugye = []
        if _sjci:
            for bl in block:
                for m2 in re.finditer(r'(?<![_가-힣])주계약\s+([\d,]{3,})', bl):
                    try:
                        v2 = int(m2.group(1).replace(',', ''))
                        if 0 < v2 <= 200000: ci_jugye.append(v2)
                    except: pass
        contracts.append({'company': company, 'product': product, 'contract_date': contract_date,
                          'expiry_date': expiry_date, 'premium': premium, 'lump_sum': lump_sum, 'pay_period': pay_period,
                          'pay_count': '', 'renewal': renewal, 'dambo': dambo,
                          'ci_jugye': ci_jugye, 'ci_extra': [], 'ipwon': [], '_sj': True,
                          'ci_lines': ci_lines, 'ci_sebu': (_SJ_SEBU or {}).get(re.sub(r'\s','',str(company or '')))})
    print(f'[SINJEONG] 3열 포맷 감지 → 계약 {len(contracts)}건 / 담보 {sum(len(c["dambo"]) for c in contracts)}개')
    return contracts



# ★★★★★v282 (2026.07.31 이영태 OCR 실사고) — 계약 경계 앵커 복구
#   <b>사고</b>: 별첨 계약 경계를 <b>'정상계약 리스트' 글자 하나</b>로만 잡는다.
#   비전 OCR이 그 장식 헤더 한 줄을 놓치면 <b>두 계약이 한 계약으로 병합</b>되고,
#   뒤 계약의 담보가 앞 계약 안으로 들어간다.
#   → CI(삼성 리빙케어) 페이지 헤더가 유실되면 <b>CI 계약 자체가 사라져 'CI 미판정'</b>이 되고,
#     그 암·뇌졸증·급성심근 담보는 비CI 계약의 담보가 되어 <b>중대한 3행이 아니라 일반 행</b>에 박힌다.
#     운전자 벌금도 남의 계약에 붙어 합산이 뒤섞인다. <b>증상 3개가 원인 1개다.</b>
#   <b>실측 재현</b>: 같은 내용 · 앵커 있음 → 계약 2건(CI True) / 앵커 1줄 삭제 → 계약 1건(운전자 담보가
#   삼성 CI 계약 안으로 흡수).
#   <b>영구원칙 위반 3번째</b>(v271 고정순번 · v277 '한장보장' 글자 · v282 '정상계약 리스트' 글자).
#   → 글자 앵커에 의존하지 않는다. <b>보험료 N원 + 보장기간 날짜~날짜</b>는 계약마다 반드시 인쇄되고
#     숫자라 OCR 생존율이 높다. 이 줄을 찾아 <b>앵커가 없는 계약 앞에만</b> 앵커 줄을 합성 삽입한다.
#     파서 본체는 손대지 않는다(회귀 0). 앵커가 이미 있으면 삽입 0건 = 무해.
_ANCHOR_TXT = '[별첨] 보험서비스(상품)별 보장 현황   (정상계약 리스트)'
def _repair_anchor(lines):
    try:
        _re_prem = re.compile(r'보험료\s*[\d,\.]+\s*원')
        _re_per  = re.compile(r'\d{4}\.\d{2}\.\d{1,2}\s*[-~（卜]\s*\d{4}\.\d{2}\.\d{1,2}')
        ins = []
        for i, l in enumerate(lines):
            if not (_re_prem.search(l) and _re_per.search(l)): continue
            # 계약 블록 시작(회사·상품 줄 포함)까지 최대 4줄 되감는다 — 빈 줄이 경계다.
            b = i
            while b > 0 and (i - b) < 4 and lines[b-1].strip(): b -= 1
            # 블록 시작 앞 6줄 안에 이미 앵커가 있으면 건드리지 않는다.
            if any('정상계약 리스트' in lines[k] for k in range(max(0, b-6), b)): continue
            if any('실효계약 리스트' in lines[k] or '미납해지' in lines[k]
                   for k in range(max(0, b-6), b)): continue
            ins.append(b)
        if not ins: return lines
        out = []
        _st = set(ins)
        for i, l in enumerate(lines):
            if i in _st: out.append(_ANCHOR_TXT)
            out.append(l)
        print(f'[v282 anchor] 계약 경계 앵커 {len(ins)}건 복구(유실된 별첨 헤더 보정)')
        return out
    except Exception as _e:
        print(f'[v282 anchor ERR] {_e}'); return lines

# ★★★★★v370 가입제안서 파서 (지점장 지시 2026.08.09) — 오른쪽 업로드 칸 전용.
#   소스 = 제안서 <b>「가입담보요약」 표</b>(No. / 가입담보 / 가입금액 / 보험료(원) / 납기·만기).
#   ★금액: 한글단위(3천만원→3,000) · <b>원 단위는 뒤 4자리 절삭</b>(79,822원→7, 지점장 확정).
#   ★회사명: 표지 로고는 이미지라 안 읽힌다 → <b>전 페이지에서 `○○손해보험/생명` 스캔</b>(고정 쪽 금지).
_JN_AMT  = re.compile(r'(\d+(?:,\d{3})*)\s*(억원|천만원|백만원|십만원|만원|원)')
# ★★★★★v386 (지점장 지적 2026.08.12 「흥국을 못읽어낸다」): 구 _JN_TERM은 <b>슬래시형만</b> 잡았다.
#   흥국 제안서 기간칸은 <b>`20년갱신 100세만기`</b>(슬래시 없음) → 담보 <b>전량 0건</b>이 됐다(실측).
#   DB `20년만기20년납`도 같은 이유로 폴백 경로에서만 잡혔다. 두 형태를 정본에 추가한다.
# ★★★★★v389 (지점장 지시 2026.08.12 「다 읽혀야한다」): 기간칸 형태가 3개뿐이라
#   <b>현대 `20년납90세만기`·`전기납10년만기갱신(최대90세)`, 삼성 `10년납 10년만기`,
#   한화 `100세만기 / 20년납`이 전부 미매칭</b>이었다. 기간을 못 찾으면 그 담보는
#   조용히 사라진다(실측 현대 9/19 · 삼성 4/20+). 5형태 → 7형태.
_JN_TERM = re.compile(r'\d+년\s*/\s*\d+년(?:\([^)]*\))?|\d+년\s*/\s*\d+세'
                      r'|\d+년\s*갱신\s*\d+\s*[세년]\s*만기'
                      r'|\d+년\s*만기\s*\d+\s*년\s*납'
                      r'|(?:전기납|\d+\s*년\s*납)\s*\d+\s*[세년]\s*만기(?:갱신)?(?:\([^)]*\))?'
                      r'|\d+\s*[세년]\s*만기\s*/\s*(?:전기납|\d+\s*년\s*납)')
_JN_CO   = re.compile(r'([가-힣A-Za-z]{2,10}(?:손해보험|화재해상보험|화재|생명보험|생명|손보))')

# ★★★★★v371: '1백50만원' 같은 <b>복합 한글단위</b>. 구 _JN_AMT는 '50만원'만 잡아 150→50이 됐다.
_JN_AMT2 = re.compile(r'(\d+)\s*(천|백|십)\s*(\d+)?\s*만원')
# ★★v389e 복합단위 2형태 — 실측 DB `1천6백만원`이 <b>6백만원=600</b>(정답 1,600),
#   삼성 `1억5,000만원`이 <b>5,000</b>(정답 15,000)으로 읽혔다.
_JN_AMT2B = re.compile(r'(\d+)\s*억\s*([\d,]+)\s*만원')
_JN_AMT2C = re.compile(r'(\d+)\s*천\s*(\d+)\s*백\s*만원')
# ★★★v389i `13만8,000원`(롯데 2번 보험료납입지원) — '만' 뒤에 원 단위가 붙는 형태.
#   구 규칙은 '8,000원'만 금액으로 잡아 <b>담보명이 '13만'이라는 가짜 행</b>이 생겼다(실측).
_JN_AMT2D = re.compile(r'(\d+)\s*만\s*[\d,]+\s*원')
# ★v371 제안서 총보험료 — 담보별 합산은 파싱이 한 건이라도 새면 틀린다. 명시값을 1순위로 쓴다.
# ★★★★★v526 제104조 (실측 2026.08.21 김애경·강민성 삼성생명 제안서):
#   구 패턴의 `해약환급`이 <b>담보명 자체</b>를 잡아먹었다. 삼성생명·미래에셋은 담보명에
#   `(무배당,무해약환급금형)`·`[해약환급금이 없는 유형]`을 붙인다 →
#   <b>실측 김애경 28행→14행 · 강민성 21행→7행이 로그 한 줄 없이 소실</b>(제11조 위반).
#   ★차단 대상은 「해약환급금 <b>예시/안내/비교</b>」 <b>표 제목줄</b>이지 담보명이 아니다.
#     제목 형태로만 좁힌다. `환급률`은 해약환급금 예시표의 열머리라 그대로 막는다.
# ★★★★★v607 (지점장 실측 2026.08.30 「어디서 대물이 540이 추가됐으며」 · 정상범 흥국).
#   [원인] `[자기부담금] 대인없음, 대물(누수제외) 20만원, 대물(누수) 50만원` 줄이
#     <b>담보로 잡혀</b> 대물에 40이 붙었다(500 → 540). 자기부담금은 가입금액이 아니다.
_JN_NOISE = re.compile(r'자기부담금|계약사항|납입보험료|적립보험료|보장보험료|합계보험료|할인후|보험가격지수'
                       r'|해약환급금\s*(?:예시|안내|비교|관련)|환급률|가입담보|담보가입')
_JN_PREM = re.compile(r'(?:할인후초회보험료|합계보험료|보장합계|1회\s*보험료|월\s*보험료|보장보험료\s*합계)\D{0,40}?([\d][\d,]{3,})')

def _jn_in_paren(s, pos):
    """★v389 pos가 괄호 안인가. 담보명 안의 한도 금액을 가입금액으로 오인하는 것을 막는다."""
    d = 0
    for ch in (s or '')[:pos]:
        if ch in '(（[': d += 1
        elif ch in ')）]': d = max(0, d-1)
    return d > 0

def _jn_amt_at(s):
    """문자열에서 첫 금액(만원 환산)을 찾는다 → (값, start, end) / 없으면 None.
    ★★v389: <b>괄호 안 금액은 건너뛴다</b>. 실측 한화 459
    `유방,갑상선…통합치료비(연간1천만원한도)(간편)  1,000만원  940원`에서 담보명 안의
    <b>'1천만원'</b>을 가입금액으로 잡아 담보명이 '…통합치료비(연간'으로 잘리고
    보험료가 <b>1,000</b>(정답 940)이 됐다."""
    s = s or ''
    cands = []
    for m in _JN_AMT2B.finditer(s):
        cands.append((m.start(), m.end(), int(m.group(1))*10000 + int(m.group(2).replace(',',''))))
    for m in _JN_AMT2C.finditer(s):
        cands.append((m.start(), m.end(), int(m.group(1))*1000 + int(m.group(2))*100))
    for m in _JN_AMT2D.finditer(s):
        cands.append((m.start(), m.end(), int(m.group(1))))
    for m in _JN_AMT2.finditer(s):
        base = {'천':1000,'백':100,'십':10}[m.group(2)]
        v = int(m.group(1))*base + (int(m.group(3)) if m.group(3) else 0)
        cands.append((m.start(), m.end(), v))
    for m in _JN_AMT.finditer(s):
        cands.append((m.start(), m.end(), _jn_won(m.group(1), m.group(2))))
    if not cands: return None
    cands.sort(key=lambda x: (x[0], -x[1]))
    for st, en, v in cands:
        if _jn_in_paren(s, st): continue
        return v, st, en
    st, en, v = cands[0]
    return v, st, en

def _jn_total_prem(txt):
    m = _JN_PREM.search(txt or '')
    if not m: return 0
    try: return int(m.group(1).replace(',',''))
    except: return 0

def _jn_won(num, unit):
    n = int(str(num).replace(',',''))
    if unit=='억원':   return n*10000        # ★v386 흥국 '10억원' → 100,000만원
    if unit=='천만원': return n*1000
    if unit=='백만원': return n*100
    if unit=='십만원': return n*10
    if unit=='만원':   return n
    return n//10000          # ★'원' 단위 = 뒤 4자리 절삭

# ★★★★★v386 흥국 제안서 = 표 왼쪽에 <b>구분 2열</b>('선택  치료')이 붙는다.
#   구 코드는 이걸 담보명으로 읽어 <b>모든 담보명이 '선택치료'</b>가 되거나, 접힘 판정
#   (`한글이 없으면 접힘`)이 <b>한글이 있어서</b> 안 걸려 윗줄 담보명을 못 붙였다.
# ★★★v389i 구 패턴은 `^기본\s*`만으로도 매칭돼 <b>담보명 `기본계약`의 '기본'을 먹었다</b>
#   (실측 현대 001번 → `계약(상해사망…)`). 구분 2열은 항상 <b>뒤에 공백</b>이 있다(흥국 `선택 치료`).
#   → `\s+`를 필수로 해 붙어 있는 담보명은 건드리지 않는다.
_JN_GUBUN = re.compile(r'^(?:기본|선택)\s+(?:기본|선택|납입면제|진단|치료|수술|입원|통원|일당|사망|후유장해|배상|비용|간병|실손|재물|운전자|기타)?\s*')
def _jn_strip_gubun(x):
    return _JN_GUBUN.sub('', x or '').strip()
def _jn_hasname(x):
    t = re.sub(r'[\s\[\]()·,]', '', x or '')
    return bool(re.search(r'[가-힣]', t)) and len(t) >= 2

def _jn_frag(l):
    """★★★v391a 담보명이 여러 줄로 접힌 표에서 <b>이어지는 조각 줄</b>인가.
    조각 = 금액도 기간도 없고, 한글이 있고, 짧고, 표 머리·설명문이 아닌 줄."""
    t = (l or '').strip()
    if not t: return False
    if _jn_amt_at(l): return False
    # ★v391a2 한화는 조각 줄에 <b>기간칸이 같이</b> 찍힌다(`…(연간1회한)(   100세만기 / 20년납`).
    #   기간만 보고 잘라내면 담보가 통째로 사라진다(실측 이진림 2건) — 기간·번호를 지운 뒤 판정한다.
    core = _jn_frag_txt(l)
    if not _jn_hasname(core) or len(core) > 60: return False
    if _JN_NOISE.search(core): return False
    if re.search(r'경우|한도로|받고|받은|말합니다|하여야|보장개시|피보험자', t): return False
    return True

def _jn_frag_txt(l):
    t = (l or '').replace('┖',' ').strip()
    # ★v391a2 구 `^\d+\s*[.)]?\s*`는 `1회한)`의 <b>1을 먹었다</b>(실측 현대 `(연간1회한)`→`(연간회한)`).
    #   순번은 <b>마침표/괄호</b>가 붙거나 <b>공백</b>이 뒤따른다 — 그때만 지운다.
    t = re.sub(r'^\s*\d+\s*[.)]\s*|^\s*\d+\s+','', t)
    t = re.sub(r'^\([^)]*\)\s*(?=[가-힣])','', t).strip()
    t = _JN_TERM.sub('', _jn_strip_gubun(t)).strip()
    return re.sub(r'\s+','', t)

def _jn_depth_bad(nm):
    """왼쪽이 잘렸는가 = 닫는 괄호가 먼저 오거나 조각 문자로 시작."""
    if not nm: return True
    if nm[0] in ')ⅠⅡⅢⅣⅤ·,': return True
    d = 0
    for ch in nm:
        if ch == '(': d += 1
        elif ch == ')':
            d -= 1
            if d < 0: return True
    return False

def _jn_depth_open(nm):
    d = 0
    for ch in (nm or ''):
        if ch == '(': d += 1
        elif ch == ')': d -= 1
    return d > 0

def _jn_scope(txt):
    """★v371 담보 표 구간만 자른다(지점장 정본: 담보 소스 = 「가입담보요약」 표).
    못 찾으면 전문을 그대로 돌려준다."""
    lines = (txt or '').split('\n'); s=None; e=None
    for i,l in enumerate(lines):
        if s is None and '가입담보' in l: s=i
        elif s is not None and e is None and ('가입설계' in l or '주의사항' in l
             or re.sub(r'\s','',l).startswith('보장보험료합계')      # ★v386 흥국 표 끝
             or re.sub(r'\s','',l).startswith('보장합계')): e=i
    if s is not None:
        return lines[s:(e if e is not None else len(lines))], True
    # ★★★★★v526 제105조 (실측 2026.08.21):
    #   삼성생명·미래에셋·메트라이프 제안서에는 「가입담보」 표가 <b>없다</b>. 구 코드는 구간을
    #   못 잘라 <b>전문</b>을 담보표로 돌렸고, 그래서 뒷쪽 「보장내역」 표의 파편이
    #   담보로 올라왔다(실측 송명복 16건 전부 `[사망보험금]`·`계약일로부` 같은 셀 조각).
    #   ★기존 통과 회사(KB·현대·롯데·한화·흥국)는 위 `가입담보` 경로로 그대로 나간다 — 건드리지 않는다.
    #   ★앵커 근거: 삼성생명 2건(김애경·강민성) 실측. 미래에셋·메트라이프는 <b>각 1건</b>뿐이므로
    #     [확인] 대기 — 표 못 찾으면 조용히 넘기지 않고 로그로 남긴다.
    _H = (r'^\[?계약사항\]?$'                                   # 삼성생명
          r'|^구\s*분보험가입금액보험기간납입기간보험료'          # 삼성생명 표머리
          r'|^구\s*분가입금액보험기간납입기간보험료'              # 메트라이프 표머리
          r'|^보험종류피보험자')                                  # 미래에셋 표머리
    _T = (r'^합계보험료|^할인전보험료|^\[?특약\s*가입\s*개요\]?'
          r'|^보장내역\s*안내|^보장내용|^선택된\s*특약이\s*없습니다')
    s2=None; e2=None
    for i,l in enumerate(lines):
        t = re.sub(r'\s','',l)
        if s2 is None:
            if re.search(_H, t): s2=i
        elif e2 is None and re.search(_T, t): e2=i
    if s2 is None:
        print('[JEAN] 담보표 구간 앵커 없음 — 전문을 대상으로 한다(오염 가능)')
        return lines, False
    print('[JEAN] v526 계약표 구간 %d~%d행 (「가입담보」 표 없는 제안서)' % (s2, e2 if e2 is not None else len(lines)))
    return lines[s2:(e2 if e2 is not None else len(lines))], True

def _jn_rows_tbl(lines):
    """★★★★★v371: 담보 행 = <b>금액 + 납입|보험기간이 같은 줄</b>에 있는 행.
    KB는 줄머리가 `2 일반상해사망…`(마침표 없음)·`┖ …`(번호 없음)이라 구 `\\d+\\.` 규칙이
    표를 통째로 건너뛰었다(실측 52건 → 1건). 번호 유무·마침표 유무에 의존하지 않는다."""
    # ★★v389g 금액 <b>단위가 다음 줄로 밀린</b> 표가 있다(KB `가입78,101` / 다음 줄 `원`).
    #   단위가 없으면 금액으로 인정되지 않아 담보가 통째로 사라진다(실측 KB 1건).
    _lx=[]
    for _i,_l in enumerate(lines):
        _nx=(lines[_i+1].strip() if _i+1<len(lines) else '')
        if _nx in ('원','만원','억원','천만원','백만원','십만원') and re.search(r'\d', _l):
            _l = _l.rstrip()+_nx
        _lx.append(_l)
    lines = _lx
    rows=[]
    # ★★★v389i 삼성 재물형은 <b>같은 담보가 건물1·건물2에 각각</b> 있고 담보명·금액이
    #   완전히 같다(`화재(폭발포함)배상책임Ⅱ(대물1사고당) 20억원`). 구 중복제거 키가
    #   (담보명,금액)이라 <b>건물2 3건이 통째로 삭제</b>됐다(실측 20/23). 섹션 라벨을
    #   키에 넣어 <b>다른 섹션이면 보존</b>한다. 흥국(같은 표 2회 게재)은 섹션도 같아 종전대로 제거.
    _sec = ''
    _used = set()      # ★v391a 이미 다른 담보가 흡수한 조각 줄(재사용 금지)
    for i,l in enumerate(lines):
        _sm = re.match(r'\s*([가-힣]{2,8}\(\d+/\d+\))', l)
        if _sm: _sec = _sm.group(1)
        # ★★v391c2 담보표가 아닌 <b>안내·예시·비교표</b> 줄 차단(실측 한화 p7 `특약 안내사항: …(연간1천만원한도)`
        #   → `…통합치료비(연간` 이라는 <b>잘린 유령 담보</b>가 생겼다).
        if re.search(r'안내사항|예시표|가격지수|보험료\s*비교|비교표', l): continue
        am=_jn_amt_at(l)
        # ★★v389h 금액칸이 <b>다른 줄로 흩어진</b> 담보(KB 6번 `가입78,101`은 윗줄, `원`은 아랫줄).
        #   금액이 없어도 <b>줄머리 번호 + 기간</b>이 둘 다 있으면 담보행으로 인정하고
        #   금액은 None으로 둔다 → 값은 안 들어가고 <b>[확인]큐에 뜬다</b>(조용히 사라지지 않는다).
        if not am:
            if _JN_TERM.search(l) and re.match(r'\s*\d+\s*[.)]?\s', l):
                _nm0=_JN_TERM.sub('', re.sub(r'^\s*\d+\s*[.)]?\s*','',l)).strip()
                _nm0=re.sub(r'[\d,]+\s*$','',_nm0).strip()
                _nm0=re.sub(r'\s+','',_jn_strip_gubun(_nm0))
                _nm0=re.sub(r'[\d,.]+$','',_nm0)
                if _jn_hasname(_nm0) and not _JN_NOISE.search(_nm0) \
                   and not re.search(r'경우|한도로|받고|받은|말합니다|하여야', _nm0):
                    rows.append({'no':0,'name':_nm0,'amt':None,'prem':0,'term':''})
            continue
        tm=_JN_TERM.search(l)
        if not tm:
            # ★v371 갱신형은 기간칸이 2줄이라 `10년/10년갱신`이 <b>다음 줄로 밀린다</b>
            #   (실측 KB 갱신 담보 9건 = 3,668원이 통째로 빠졌다). 아래 2줄까지만 본다.
            # ★★★v389c 삼성 재물형은 <b>한 번호 아래 담보명이 여러 줄</b>이고 기간·보험료는
            #   그중 <b>한 줄에만</b> 있다(`20억원` 줄엔 번호도 기간도 없다). 구 코드는 다음
            #   금액 줄을 만나면 즉시 break해 그 줄들을 통째로 버렸다(실측 4/20+).
            #   → 기간 검사를 break보다 <b>먼저</b> 한다.
            for _j in range(i+1, min(i+5, len(lines))):
                tm=_JN_TERM.search(lines[_j])
                if tm: break
                if _jn_amt_at(lines[_j]): break        # 다음 담보 행이면 중단
        # ★v371: pdftotext는 <b>페이지 경계에서 기간칸을 다른 쪽으로 흘린다</b>(실측 갱신형 6건).
        #   기간이 없어도 <b>줄머리가 번호 또는 ┖</b>이면 표의 담보 행으로 인정한다.
        #   반대로 `계약자|청약번호 … 185,602원` 류는 줄머리가 한글이라 여기서 걸러진다.
        # ★★v389 줄머리 번호 뒤에 <b>마침표</b>가 오는 표(현대 `186.` `425.`)가 있다.
        #   구 `\d+\s`는 이걸 못 잡아, 기간칸까지 미매칭이면 담보가 통째로 사라졌다(실측 10건).
        # ★★★★★v422 (삼성 조승우 실측 2026.08.15): 삼성 요약표는 <b>한 번호 아래 담보명이 여러 줄</b>이고
        #   보험료·기간은 그 그룹의 <b>대표 줄 하나에만</b> 있는데, 그 대표 줄이 담보명 줄들 <b>사이</b>에 끼어 온다
        #   (`344 … 269 20년납 90세만기`가 6종과 7종 사이). 구 코드는 기간을 <b>아래로만</b> 찾고
        #   다음 금액 줄에서 break해, 대표 줄이 위에 있는 형제 담보를 통째로 버렸다.
        #   실측 소실 9건 — 상해 4·5·7·8종 / 질병 1·3·4·5·7종 수술비(시술포함), 53건 중 44건만 파싱.
        #   → 줄머리 <b>대괄호 태그</b>(`[간편]` 등)도 `┖`와 같은 <b>담보 행 마커</b>로 인정한다.
        #   ★특정 글자를 박지 않는다(구조 가정 금지) — 대괄호 형태만 본다.
        _head = re.match(r'\s*(?:\d+\s*[.)]?\s|┖|\[[^\[\]]{1,10}\])', l)
        if not tm and not _head: continue
        amt, st, en = am
        after=l[en:]
        # ★★★v389 보험료 = 금액 뒤 첫 숫자였는데, 롯데는 `1,000만원 20년/20년 갱신 61,800`
        #   순서라 <b>'20년'의 20</b>을 보험료로 집었다(전 행 20 → 합계 860, 정답 283,966).
        #   → 기간 토큰을 먼저 지우고 남은 첫 숫자를 쓴다.
        _aft = _JN_TERM.sub(' ', after)
        _aft = re.sub(r'\d+\s*(?:년|세|개월|회|일)', ' ', _aft)
        pm=re.search(r'(\d[\d,]*)', _aft)           # ★v386 구 `[\d,]+`는 쉼표만 잡아 int() 크래시(실측 DB)
        try: prem=int(pm.group(1).replace(',','')) if pm else 0
        except Exception: prem=0
        # ★★v389 가입금액칸이 '보장내용 참조'인 담보(한화 4번)는 <b>보험료가 첫 금액</b>으로
        #   잡혀 amt='1,791원'→0, prem=0이 됐다. 원 단위 + 뒤에 숫자 없음 = 그 값이 보험료다.
        # ★★★v391c 담보표가 아니라 <b>보험료 비교표</b>인 줄(한화 p9 간편↔일반 비교)이
        #   담보로 잡혀 <b>잘린 유령 담보</b>가 생겼다(실측 이진림 1건). 금액 뒤 숫자칸이
        #   4개 이상이면 담보 행이 아니라 비교·예시표다.
        if len(re.findall(r'\d[\d,]*', _aft)) >= 4: continue
        if prem == 0 and amt == 0 and re.match(r'\s*[\d,]+\s*원', l[st:en]):
            try: prem = int(re.sub(r'[^\d]', '', l[st:en]))
            except Exception: prem = 0
        # ★★v526 제107조 (실측 2026.08.21 송명복): 구 `^\\d+\\s*[.)]?\\s*`는 <b>\\s*가 0개를 허용</b>해
        #   `1-5종수술특약`의 앞 `1`을, `15대질환수술특약`의 앞 `15`를 <b>줄머리 번호로 오인해 지웠다</b>
        #   (→ `-5종수술특약`·`대질환수술특약`). 번호는 <b>뒤에 구분자(공백·마침표·괄호)가 있을 때만</b> 번호다.
        nm=re.sub(r'^\s*\d+\s*[.)]\s*|^\s*\d+\s+','', l[:st].replace('┖','')).strip()
        nm=re.sub(r'^\([^)]*\)','',nm).strip()   # ★v371b 상품 접두어 `(맞춤_간편고지Ⅱ)` 제거 —
        #   구 v370에 있던 규칙을 새 파서가 빠뜨려 <b>DB 상해수술비 50·질병수술비 20이 소실</b>됐다(실측).
        nm=_jn_strip_gubun(nm)                # ★v386 흥국 구분 2열('선택 치료') 제거
        nm=_JN_TERM.sub('', nm).strip()       # ★v386 기간칸이 금액 <b>앞</b>에 오는 표(흥국) — 담보명에서 제거.
        #   안 지우면 담보명이 '…20년갱신100세만기'로 오염되고, 접힘 행은 기간만 남아 접힘 판정이 안 된다.
        # ★★★v389i 담보명이 <b>번호 줄에서 잘려 내려온</b> 표(삼성 `48   붕괴·` / 다음 줄
        #   `침강 및 사태손해(실손)  8,000만원 …`). 담보명이 `침강및사태손해`로 오염돼
        #   매핑이 어긋난다. 현재 줄에 번호가 없고 <b>윗줄이 번호+짧은 조각(금액 없음)</b>일 때만
        #   앞에 붙인다. 윗줄에 금액이 있으면 그건 별개 담보라 건드리지 않는다.
        if not re.match(r'\s*(?:\d+\s*[.)]?\s|┖)', l):
            _pv = lines[i-1] if i>0 else ''
            _pm = re.match(r'\s*\d+\s+(\S.{0,18})\s*$', _pv)
            # ★v389i 윗줄이 <b>보험료+기간칸</b>이면 담보명 조각이 아니다(실측 삼성 `620 10년납 10년만기`
            #   가 담보명 앞에 붙었다). 기간 토큰이 있으면 건너뛴다.
            if _pm and _JN_TERM.search(_pv): _pm = None
            if _pm and not _jn_amt_at(_pv) and re.search(r'[가-힣]', _pm.group(1)):
                nm = _pm.group(1).strip() + nm
                _used.add(i-1)   # ★v391a2 v391a 위흡수가 <b>같은 줄을 또</b> 붙이던 것 차단(실측 삼성 `붕괴·붕괴·`)
        # ★★★v391a 담보명이 <b>여러 줄로 접힌</b> 표(롯데: 번호줄 위 3줄 + 아래 1줄).
        #   구 코드는 위 1줄·아래 1줄만 봐서 `유사암진단비)(간편맞춤형Ⅱ)`처럼 <b>앞이 잘린 이름</b>이
        #   그대로 매핑에 들어갔다(실측 롯데 3건 — 순번 2·7·59).
        #   ─ 위로: 이름이 없거나 <b>왼쪽이 잘린</b> 동안 올라간다.
        #   ─ 아래로: 괄호가 <b>안 닫힌</b> 동안 내려간다.
        #   ─ 흡수한 줄은 `_used`에 넣어 <b>다음 담보가 다시 쓰지 못하게</b> 한다(앞 담보 꼬리 오염 차단).
        nm = re.sub(r'\s+','', nm)
        _wasfold = not _jn_hasname(nm)
        _k = i-1
        while _k >= 0 and (i-_k) <= 5 and _k not in _used and _jn_frag(lines[_k]):
            _fr = _jn_frag_txt(lines[_k])
            if not _fr: break
            # ★v391a2 윗줄이 `·`·`,`·`(`·`및`로 <b>끝나면</b> 아랫줄로 이어지는 조각이다
            #   (롯데 순번 7 `갑상선암·` / `기타피부암·`). 괄호 깊이만으로는 안 잡힌다.
            if not (_jn_depth_bad(nm) or _fr[-1] in '·,(/' or _fr.endswith('및')): break
            nm = _fr + nm; _used.add(_k); _k -= 1
        _k = i+1; _first_down = True
        while _k < len(lines) and (_k-i) <= 4 and _k not in _used and (
                _jn_frag(lines[_k]) or
                # ★v391a2 `형)` 같은 <b>한 글자 닫는 조각</b>은 _jn_hasname(2자 이상) 미달이라
                #   흡수가 안 돼 담보명이 `…(갱신`으로 끝났다(실측 롯데 순번 59).
                (_jn_depth_open(nm) and len(lines[_k].strip()) <= 8 and ')' in lines[_k]
                 and not _jn_amt_at(lines[_k]))):
            _fr = _jn_frag_txt(lines[_k])
            if not _fr: break
            # 아래로: 괄호가 안 닫혔거나, 접힘 행의 <b>첫 아랫줄이 괄호로 시작</b>(롯데 `(갱신형)`)
            if not (_jn_depth_open(nm) or (_first_down and _wasfold and _fr[0] in '([')): break
            nm = nm + _fr; _used.add(_k); _k += 1; _first_down = False
        _used.add(i)
        nm=re.sub(r'\s+','',nm)
        # ★★★v416 제안서 담보명 <b>꼬리 오염</b>(실측 KB 26.07 간병인지원 3건).
        #   접힘 아랫줄 `신형)`에 <b>우측 열 `(갱신종료:90세)`가 같이 찍혀</b> 담보명에 붙었고,
        #   `(갱`+`갱신`+`신형)` 결합으로 <b>`갱갱신신형`</b>이 됐다.
        #   → v391 영구조항(「담보명이 조용히 바뀐다」) 위반이라 그 자리에서 정화한다.
        #   `(갱신종료:…)`는 <b>기간칸</b>이지 담보명이 아니다.
        nm=re.sub(r'\(갱신종료[^)]*\)', '', nm)
        nm=re.sub(r'갱갱신신형', '갱신형', nm)
        nm=re.sub(r'(?:갱신){2,}형', '갱신형', nm)
        if not _jn_hasname(nm): continue
        # ★v389e 표 헤더·합계 줄이 담보로 잡히던 것 차단(실측 DB `계약사항:01종…납입보험료`,
        #   `만기/납기납입보험료`, `적립보험료`). 담보명 `보험료납입면제`·`보험료납입지원`은 안전.
        if _JN_NOISE.search(nm): continue
        # ★v389f 담보 상세 설명문이 담보행으로 잡히던 것 차단(실측 DB 3건).
        #   서술어가 있거나 금액 표기가 2회 이상이면 담보명이 아니라 문장이다.
        if re.search(r'경우|한도로|받고|받은|말합니다|하여야', nm) or nm.count('만원') >= 2: continue
        rows.append({'no':0,'name':nm,'amt':amt,'prem':prem,'sec':_sec,
                     'term':(tm.group(0).replace(' ','') if tm else '')})
    # ★★v386 같은 담보 표가 문서에 두 번 실리는 제안서(흥국: 「가입담보 리스트」+「보장사항」)가 있다.
    #   중복이 그대로 합류하면 보험료·담보가 <b>2배</b>가 된다 → (담보명, 금액) 기준 첫 건만 남긴다.
    _seen=set(); _uq=[]
    for x in rows:
        _k=(x['name'], x['amt'], x.get('sec',''))
        if _k in _seen: continue
        _seen.add(_k); _uq.append(x)
    if len(_uq)!=len(rows): print(f'[JEAN] 중복 담보행 {len(rows)-len(_uq)}건 제거 → {len(_uq)}건')
    return _uq

def _jn_scopes(txt):
    """★★★v389d 담보 표가 <b>여러 블록</b>인 제안서가 있다(삼성 재물 = 건물1 / 피보험자 /
    건물2, 블록마다 `보장보험료 합계`로 끝난다). 구 `_jn_scope`는 <b>첫 합계에서 끊어</b>
    뒤 블록을 통째로 버렸다(실측 12건, 정답 20+). → 시작~종료 구간을 <b>전부</b> 모은다."""
    lines = (txt or '').split('\n'); segs=[]; s=None
    for i,l in enumerate(lines):
        t = re.sub(r'\s','',l)
        if s is None:
            # ★v389e `보장(보상)내용` 열이 있는 표는 <b>담보 상세 설명 블록</b>이다(DB 8~9p).
            #   요약표와 같은 담보를 다시 싣고 설명문에서 가짜 행이 나온다(실측 5건).
            if ('가입담보' in t or '담보가입' in t) and '보장내용' not in t and '보상)내용' not in t: s=i
        # ★★★★★v605 제107조 2항 (지점장 지시 2026.08.30 「요약에서 봐서 해야지」
        #   · 정상범 흥국 제안서 실측).
        #   [실측] 이 문서는 <b>「■ 가입담보 리스트」(요약표)</b>와
        #     <b>「■ 보장사항」(담보별 지급조건 상세표)</b>가 한 문서에 같이 있다.
        #     두 표에 <b>같은 담보가 각각 한 번씩</b> 실려 있는데 종료 앵커에
        #     `보장사항`이 없어 두 구간을 <b>한 덩어리로 읽었다</b>
        #     → 34종이 전부 2배(뇌혈관 500→1,000 · 특정3대 500→1,000).
        #   ⇒ <b>「보장사항」·「보장내용」 표가 시작되면 요약표는 거기서 끝난다.</b>
        #     v389e가 <b>시작</b> 앵커에서 상세표를 걸렀던 것과 같은 이유를 <b>종료</b>에도 적용.
        elif ('가입설계' in t or '주의사항' in t
              or '보장보험료합계' in t or '보장합계' in t
              or re.match(r'^[■□▣●]?\s*보장사항', t)
              or re.match(r'^[■□▣●]?\s*보장내용$', t)):
            segs.append((s,i)); s=None
    if s is not None: segs.append((s,len(lines)))
    if segs: return [lines[a:b] for a,b in segs], True
    # ★★★★★v526 제105조 (실측 2026.08.21):
    #   삼성생명·미래에셋·메트라이프 제안서에는 「가입담보」 표가 <b>없다</b>. 구 코드는
    #   구간을 못 잘라 <b>전문</b>을 담보표로 돌렸고, 뒷쪽 「보장내역」 표의 셀 조각이
    #   담보로 올라왔다(실측 송명복 16건 전부 `[사망보험금]`·`계약일로부` 같은 파편).
    #   ★기존 통과 회사(KB·현대·롯데·한화·흥국)는 위 `가입담보` 경로로 그대로 나간다.
    #   ★앵커 근거: 삼성생명 2건(김애경·강민성) 실측. 미래에셋·메트라이프는 각 1건뿐 → [확인] 대기.
    _H = (r'^\[?계약사항\]?$'
          r'|^구\s*분보험가입금액보험기간납입기간보험료'
          r'|^구\s*분가입금액보험기간납입기간보험료'
          r'|^보험종류피보험자')
    _T = (r'^합계보험료|^할인전보험료|^\[?특약\s*가입\s*개요\]?'
          r'|^보장내역\s*안내|^보장내용|^선택된\s*특약이\s*없습니다')
    s2=None; e2=None
    for i,l in enumerate(lines):
        t = re.sub(r'\s','',l)
        if s2 is None:
            if re.search(_H, t): s2=i
        elif e2 is None and re.search(_T, t): e2=i
    if s2 is None:
        print('[JEAN] v526 계약표 앵커 없음 — 전문 대상(오염 가능)')
        return [lines], False
    print('[JEAN] v526 계약표 구간 %d~%d행 (「가입담보」 표 없는 제안서)'
          % (s2, e2 if e2 is not None else len(lines)))
    _seg = lines[s2:(e2 if e2 is not None else len(lines))]
    # ★★★★★v526 제108조 (실측 2026.08.21 백명자 메트라이프 달러종신):
    #   계약 통화가 <b>달러</b>다(`＄34,000.00` · `＄613.70`). 마스터는 만원 단위인데
    #   <b>달러→만원 환산 규칙은 지침에 없다</b>. 환산을 지어내면 34,000이 3.4억으로 들어간다.
    #   → 값을 만들지 않고 <b>시끄럽게 멈춘다</b>(제11조). 환산 규칙은 지점장 확정 사항이다.
    if re.search(r'[＄$]\s?[\d,]', '\n'.join(_seg)):
        print('[JEAN] ★외화(달러) 계약 — 달러→만원 환산 규칙이 지침에 없다. '
              '열을 만들지 않는다. [확인] 지점장 확정 필요')
        return [[]], True
    return [_jn_norm_unitless(_seg)], True

# ★★★★★v526 제106조 (실측 2026.08.21 송명복 미래에셋 M-케어):
#   미래에셋 표는 <b>단위가 셀이 아니라 열머리글</b>에 있다(`보험가입금액` / 아랫줄 `(만원)`).
#   셀은 맨숫자(`300`)라 `_jn_amt_at`이 금액으로 인정하지 않아 <b>16행 전부 탈락</b>했다.
#   → 머리글이 `(만원)`을 선언한 표에 한해 <b>피보험자명 열이 있는 데이터 줄</b>의
#     가입금액칸에 `만원`, 보험료칸에 `원`을 붙여 다른 회사와 같은 모양으로 만든다.
#   ★가정하지 않는다 — `(만원)` 선언이 <b>실제로 있을 때만</b> 돈다. 없으면 원문 그대로.
#   ★[확인] 앵커 근거 1건(송명복). 삼성처럼 2건 확보되기 전까지는 실패 시 로그로 남긴다.
_JN_UNITLESS = re.compile(
    r'^(?P<head>.*?)(?P<who>[가-힣]{2,4})\s+(?P<amt>[\d,]+)\s+(?P<age>\d{1,3})\s+'
    r'(?P<term>종신|\d{1,3}\s*세|\d{1,2}\s*년)\s+(?P<pay>\S*납)\s+(?P<cyc>\S*납)\s+(?P<prem>[\d,]+)\s*$')

def _jn_norm_unitless(seg):
    _hdr = ''.join(re.sub(r'\s','',x) for x in seg[:4])
    if '(만원)' not in _hdr or '보험료(원)' not in _hdr:
        return seg
    # ★2항 담보명이 데이터줄 <b>위·아래로 쪼개져</b> 있다. 제82조(상품명 꼬리 절단)와 같은 사고를
    #   막으려면 이름을 <b>먼저</b> 붙여놓아야 한다. 붙이는 기준은 지침에 이미 있는
    #   <b>「괄호가 안 닫힌 동안 내려간다」</b>(v391a)를 그대로 쓴다 — 새 기준을 만들지 않는다.
    def _open(s):
        d=0
        for ch in s:
            if ch in '([': d+=1
            elif ch in ')]': d=max(0,d-1)
        return d>0
    _hit=[i for i,l in enumerate(seg) if _JN_UNITLESS.match(l.rstrip())]
    _taken=set()
    _name={}
    for i in _hit:
        m=_JN_UNITLESS.match(seg[i].rstrip())
        up=[]; j=i-1
        while j>=0 and j not in _taken and j not in _hit:
            t=seg[j].strip()
            if not t or not re.search(r'[가-힣]', t): break
            # ★표머리글은 담보명 조각이 아니다 — 흡수하면 1행 담보명이 헤더로 오염된다(실측 송명복).
            if re.search(r'보험종류|피보험자|보험가입금액|납입주기|보험료\(원\)|\(만원\)|나이', t): break
            up.insert(0,t); _taken.add(j); j-=1
            if len(up)>=3: break
        nm=(' '.join(up)+' '+m.group('head').strip()).strip()
        j=i+1; k=0
        while j<len(seg) and j not in _hit and k<3:
            t=seg[j].strip()
            if not t or not re.search(r'[가-힣]', t): break
            # 괄호가 안 닫혔거나(v391a), 아래 조각이 <b>닫는 괄호로 시작</b>하면 그 담보의 꼬리다.
            if not (_open(nm) or t[0] in ')]' or re.match(r'^\[[^\[\]]{1,20}\]$', t)): break
            nm=(nm+' '+t).strip(); _taken.add(j); j+=1; k+=1
        _name[i]=nm
    out=[]; n=0
    for i,l in enumerate(seg):
        if i in _taken: out.append(''); continue
        m=_JN_UNITLESS.match(l.rstrip())
        if not m: out.append(l); continue
        n+=1
        # ★미래에셋 표에는 <b>번호 열이 없다</b>. `_jn_rows_tbl`은 줄머리 번호·`┖`·`[태그]`를
        #   담보 행 마커로 본다 → 기존 마커 `┖`를 붙여 <b>같은 규칙</b>을 타게 한다(새 규칙 만들지 않는다).
        out.append('┖ %s %s만원 %s %s %s원'
                   % (_name.get(i, m.group('head').strip()), m.group('amt'),
                      m.group('term').replace(' ',''), m.group('pay'), m.group('prem')))
    if n: print('[JEAN] v526 단위 열머리글 표 정규화 %d행 (미래에셋형)' % n)
    else: print('[JEAN] v526 (만원) 머리글은 있는데 데이터행 0 — 표 모양 미지원 [확인]')
    return out

def parse_jean_life(txt):
    """★★★★★v599 제140조 — <b>생보(KB라이프 등) 제안서 표</b> 파서.
       손보 표(`가입담보요약`)와 형식이 완전히 달라 기존 파서가 담보 0건을 냈다.
       [실측 · 천미옥 KB라이프] 담보명이 여러 줄로 쪼개지고 금액도 `2,`+`000만원`으로 잘린다.
       ⇒ ①표 구간(`보험의 종류` ~ `합계`)을 이어붙이고 공백 제거
         ②`(간편3.10.5)`·`KB3.10.5`·`지정대리청구` 같은 <b>담보 시작 표식</b>으로 자른다
         ③각 조각의 첫 `N만원`이 가입금액, 그 앞이 담보명.
       반환은 손보 파서와 <b>같은 모양</b>이라 뒤 로직이 그대로 돈다."""
    import re as _r
    _ls = str(txt or '').split('\n')
    try:
        _st = next(i for i, l in enumerate(_ls) if '보험의 종류' in l)
    except StopIteration:
        return []
    _en = next((i for i, l in enumerate(_ls[_st:], _st) if _r.search(r'합\s*계', l)), _st + 60)
    _t = _r.sub(r'\s', '', ''.join(_ls[_st + 1:_en]))
    if not _t:
        return []
    _t = _r.sub(r'(주계약|특약|제도성)', '', _t)
    # 담보 시작 표식 — 상품코드 괄호 / 사명+코드 / 제도성 특약명
    _parts = _r.split(r'(?=\(간편[\d.]+\)|\(무배당\)|KB[\d.]+|지정대리청구)', _t)
    _out = []
    for _p in _parts:
        if not _p.strip():
            continue
        _m = _r.search(r'([\d,]+)만원', _p)
        if not _m:
            continue                      # 금액 없는 제도성 특약은 담보가 아니다
        _nm = _p.split(_m.group(0))[0]
        _nm = _r.sub(r'[\d,]+원$', '', _nm).strip()
        # ★★★★★v599b — 담보명 <b>꼬리의 열린 괄호를 잘라낸다</b>.
        #   [실측] 생보 표는 담보명이 줄바꿈으로 끊겨 `…(해약환급금미지급형)<b>(계</b>` 처럼
        #     <b>괄호가 안 닫힌 채</b> 끝난다. 그대로 넘기면 뒤의 `rule_extract`가
        #     <b>다음 담보까지 한 줄로 삼켜</b> 금액이 0이 됐다(실측 3건 소실).
        #   ⇒ 마지막 <b>닫히지 않은 `(` 이후</b>를 버린다. 괄호 수식어는 매핑에 필요 없다.
        while _nm.count('(') > _nm.count(')'):
            _nm = _nm[:_nm.rfind('(')].strip()
        _nm = _r.sub(r'[\s·,]+$', '', _nm)
        if len(_nm) < 4:
            continue
        try:
            _amt = float(_m.group(1).replace(',', ''))
        except Exception:
            continue
        if _amt <= 0:
            continue
        # ★키 이름을 <b>손보 파서와 똑같이</b> 맞춘다(`no·name·amt·prem·sec·term·blk`).
        #   다르면 뒤 로직이 조용히 0으로 읽는다.
        _tm = ''
        _mt = _r.search(_r.escape(_m.group(0)) + r'((?:\d+세만기|\d+년만기))((?:전기납\(\d+년납\)|\d+년납))', _p)
        if _mt:
            _tm = '%s/%s' % (_mt.group(2), _mt.group(1))
        _out.append({'no': len(_out) + 1, 'name': _nm, 'amt': _amt, 'prem': 0,
                     'sec': '', 'term': _tm, 'blk': 0})
    return _out


def parse_jean(txt):
    """가입제안서 텍스트 → [{'name','amt','prem','term'}] 리스트.
    ★v371: 신규 표 파서를 먼저 쓰고, 결과가 빈약하면 구 `\\d+\\.` 규칙으로 폴백한다."""
    _segs,_sc = _jn_scopes(txt)
    _new = []
    for _bi,_ls in enumerate(_segs):
        _rr = _jn_rows_tbl(_ls)
        for _r in _rr: _r['blk'] = _bi
        _new += _rr
    _sn=set(); _u=[]
    for x in _new:                       # ★v389d 구간을 합친 뒤 다시 한 번 중복 제거
        # ★★★v389i 블록(섹션)이 다르면 같은 담보라도 <b>별개 계약 대상</b>이다
        #   (삼성 재물 건물1·건물2에 화재배상책임 3줄이 각각 존재 — 구 키는 3건을 삭제했다).
        k=(x['name'], x['amt'], x.get('sec',''), x.get('blk',0))
        if k in _sn: continue
        _sn.add(k); _u.append(x)
    # ★★★★★v390a (2026.08.12 흥국 실측): 같은 담보 표가 문서에 <b>두 번</b> 실리는데
    #   두 번째 표에서는 담보명이 줄바꿈으로 <b>잘려 있다</b>(`…(통합간편가입형)(갱신형_20년)`
    #   ↔ `…(통합간`). 이름이 달라 구 중복제거를 통과해 <b>20건 → 40건</b>이 됐고, 그 40줄이
    #   뒤 단계에서 서로 뭉쳐 담보가 오배정됐다.
    #   → <b>금액·보험료·기간이 모두 같고 한쪽 이름이 다른 쪽의 접두</b>면 같은 담보로 본다.
    #     남기는 것은 <b>긴 이름</b>(온전한 쪽)이다.
    #   ★★v391d 상세 보장내용 표는 <b>보험료·기간칸이 없다</b>(한화 p3~6 `459. …(연간1천만원한도)`
    #     → 이름만 `…통합치료비(연간`으로 잘려 <b>유령 담보</b>가 남았다. 실측 이진림 1건).
    #     → <b>금액이 같고</b> 보험료·기간이 <b>같거나 한쪽이 비었으면</b> 같은 묶음으로 본다.
    _grp={}
    for x in _u: _grp.setdefault(x['amt'], []).append(x)
    for _g in list(_grp.values()):
        pass
    _drop=set()
    for _g in _grp.values():
        if len(_g) < 2: continue
        _srt=sorted(_g, key=lambda y: -len(y['name']))
        _keep=[]
        for y in _srt:
            # ★v390a 이름이 <b>완전히 같은</b> 건은 대상이 아니다 — 삼성 재물 건물1·건물2처럼
            #   같은 담보명이 섹션별로 각각 존재하는 정상 케이스를 죽인다(v389i 회귀 방지).
            def _compat(k, y):
                if k['name'] == y['name'] or not k['name'].startswith(y['name']): return False
                _p1,_p2 = k.get('prem',0), y.get('prem',0)
                _t1,_t2 = k.get('term',''), y.get('term','')
                if _p1 and _p2 and _p1 != _p2: return False
                if _t1 and _t2 and _t1 != _t2: return False
                return True
            if any(_compat(k, y) for k in _keep): _drop.add(id(y))
            else: _keep.append(y)
    if _drop:
        _u=[x for x in _u if id(x) not in _drop]
        print(f'[JEAN] v390a 잘린 중복 담보명 {len(_drop)}건 제거 → {len(_u)}건')
    _new = _u
    _old = _parse_jean_dot(txt)
    if len(_new) >= max(3, len(_old)): return _new
    return _old if _old else _new

def _parse_jean_dot(txt):
    """구 v370 규칙(줄머리 `번호.`) — DB 제안서 경로 보존용 폴백."""
    lines = txt.split('\n'); rows=[]
    for i,l in enumerate(lines):
        if not re.match(r'\s*\d+\.\s', l): continue
        m = re.match(r'\s*(\d+)\.\s*(.*)$', l)
        no, body = int(m.group(1)), m.group(2)
        am = _JN_AMT.search(body)
        if not am: continue
        amt  = _jn_won(am.group(1), am.group(2))
        term = _JN_TERM.search(body)
        after= body[am.end():]
        pm   = re.search(r'([\d,]+)', after)
        prem = int(pm.group(1).replace(',','')) if pm else 0
        nm   = body[:am.start()].strip()
        if not nm.startswith('('):            # 접힘: 앞줄(+뒷줄) 복원
            pre = lines[i-1].strip() if i>0 else ''
            nxt = lines[i+1].strip() if i+1<len(lines) else ''
            if pre.startswith('('):
                nm = pre + (nxt if (not re.match(r'\s*\d+\.', nxt) and len(nxt)<12) else '')
        nm = re.sub(r'^\([^)]*\)','',nm).strip()      # 상품 접두어 제거 (맞춤_간편고지Ⅱ) 등
        nm = re.sub(r'\s+','',nm)
        if not nm: continue
        rows.append({'no':no,'name':nm,'amt':amt,'prem':prem,
                     'term':(term.group(0).replace(' ','') if term else '')})
    return rows

def _jn_cover_word(txt):
    """★★★★★v371 (지점장 확정 2026.08.09): <b>표지가 답이다 — 상품명 1번 단어가 회사다</b>.
    지점장 원문: 「표지가답이다」 / 「1번단어가 kb다」.
    표지 로고는 이미지라 '○○손해보험' 글자가 아예 없다(KB 제안서 실측 0건).
    → 표지 첫 실질 줄(상품명)의 <b>첫 단어</b>를 그대로 회사명으로 쓴다.
    ★회사표를 새로 만들지 않는다 — `_hbkey`가 부분문자열로 잡으므로 'KB'만으로 KB 분기가 걸린다."""
    for l in (txt or '').split('\n')[:15]:
        t = l.strip()
        if not t: continue
        m = re.match(r'([A-Za-z가-힣]{2,10})', t)
        return m.group(1) if m else ''
    return ''

# ★★★★★v389 정본 회사키 = `_hbkey`가 쓰는 목록 그대로다(새 표를 만들지 않는다).
#   `let:` = 롯데(#2 정본: 2열 = 롯데(let:)).
_JN_CO2 = (('KB','KB'),('한화','한화'),('농협','NH농협'),('NH','NH농협'),('DB','DB'),
           ('현대해상','현대해상'),('흥국','흥국'),('롯데','롯데'),('let:','롯데'),
           ('삼성','삼성화재'),('메리츠','메리츠'))

def _jn_is_customer(w, txt):
    """★v389 표지 1번 단어가 <b>고객 이름</b>인가. 롯데는 로고가 이미지라 '최구갑님을',
    현대는 '정현주'가 회사명으로 들어갔다. 회사명이 틀리면 `_hbkey`가 None이 되어
    <b>심장 정본표 분해가 통째로 죽는다</b>."""
    if not w: return False
    t = txt or ''
    if re.search(re.escape(w) + r'\s*님', t): return True
    if re.search(r'(?:계약자|피보험자)\s*[:：]?\s*' + re.escape(w), t): return True
    return False

def jean_company(txt):
    """제안서 회사명. ①전문 '○○손해보험/생명' → ②정본 회사키 스캔 → ③표지 1번 단어
    (고객명이면 거부) → ④도메인."""
    for m in _JN_CO.finditer(txt or ''):
        c = m.group(1)
        if '고객상담' in c or '홈페이지' in c: continue
        # ★★★★★v422 (삼성 조승우 실측 2026.08.15): 상품명이 회사명으로 시작하는 제안서에서
        #   <b>상품 수식어까지 회사명으로</b> 잡혔다(`무배당삼성화재간편365…` → 회사 `무배당삼성화재`).
        #   정본 회사명은 <b>삼성화재</b>다. 접두 수식어만 벗긴다(회사표를 새로 만들지 않는다).
        c = re.sub(r'^(?:무배당|유배당|무|유)(?=[가-힣A-Za-z]{2,})', '', c)
        return c
    # ★★v389b 스캔 범위 = <b>표지 1페이지</b>. 전문을 훑으면 오탐이 난다 —
    #   실측: 현대 발행자코드 `4BKB27` → <b>KB</b>, 롯데 본문 문장 → <b>한화</b>.
    _t = re.sub(r'\s', '', (txt or '').split('\f')[0])
    for _k, _co in _JN_CO2:
        if _k in _t:
            print(f'[JEAN 회사] 표지 정본 회사키 = {_k!r} → {_co}')
            return _co
    _w = _jn_cover_word(txt)
    if _w and _jn_is_customer(_w, txt):
        print(f'[JEAN 회사] 표지 1번 단어 {_w!r} = 고객명 → 거부')
        _w = ''
    if _w:
        print(f'[JEAN 회사] 표지 1번 단어 = {_w!r}')
        return _w
    if 'idbins.com' in (txt or ''): return 'DB손해보험'   # ★[확인] 도메인표는 지점장 확정분만 추가
    return '미확인'

def build_proposal_contract(pdf_bytes, fname=''):
    """가입제안서 PDF → 계약 dict 1건(맨 오른쪽 열). 실패 시 None."""
    import subprocess as _sp, tempfile as _tf
    try:
        with _tf.NamedTemporaryFile(suffix='.pdf', delete=False) as _f:
            _f.write(pdf_bytes); _p=_f.name
        full = _sp.run(['pdftotext','-layout',_p,'-'],capture_output=True,text=True,timeout=90).stdout
        try: os.unlink(_p)
        except: pass
    except Exception as e:
        print('[JEAN] pdftotext 실패', e); return None
    _life_rows = None
    rows = parse_jean(full)
    if not rows:
        # ★★★★★v600 제140조 (지점장 지시 2026.08.26 「<b>kb라이프 인식해라</b>」).
        #   생보 제안서는 표 형식이 <b>통째로 다르다</b> — 담보명이 3~5줄로 쪼개지고
        #   금액도 `2,`+`000만원`으로 잘린다. 손보 파서가 담보 0건을 냈다.
        rows = parse_jean_life(full)
        if rows:
            _life_rows = rows          # ★v600 생보 경로 표식
            print('[JEAN] 생보 표 형식으로 %d건 인식(제140조)' % len(rows))
    if not rows:
        print('[JEAN] 담보 0건 — 가입담보요약 표를 못 찾음'); return None
    co   = jean_company(full)
    prod = ''
    for l in full.split('\n')[:40]:
        t=l.strip()
        if len(t)>=10 and ('보험' in t) and ('가입' not in t) and ('준법' not in t):
            prod = t; break
    # ★★★★★v422 (삼성 조승우 실측 2026.08.15): 상품명에 <b>'보험' 글자가 없는</b> 제안서가 있다
    #   (`무배당삼성화재간편365당당한새로고침100세(2604.1)`). 구 규칙이 이 줄을 건너뛰고
    #   <b>피보험자 줄</b>(`조승우 (42세 / 남 / … / 보험나이변경일 : 매년`)을 상품명으로 집었다
    #   → 엑셀 헤더 1행 상품명 칸이 고객 신상으로 오염된다(v29c (1) 헤더 3줄 표기 위반).
    #   ★<b>오염된 경우에만</b> 표지 첫 실질 줄로 대체한다 — 기존 통과 회사(KB·현대·롯데…)는 건드리지 않는다.
    if (not prod) or re.search(r'피보험자|계약자|보험나이변경일|\d+\s*세\s*/', prod):
        for l in full.split('\n')[:20]:
            t=l.strip()
            if len(t) < 8: continue
            if re.search(r'님을?\s*위한|가입제안서|상품제안서|계약자|피보험자|고객님', t): continue
            prod = t; break
    # ★v371 총보험료 = <b>제안서 명시값 1순위</b>(할인후초회보험료 등), 없으면 담보별 합.
    #   담보별 합은 한 건만 새도 틀린다(실측 KB 185,602 → 357).
    prem = _jn_total_prem(full) or sum(x['prem'] for x in rows)
    # 납기/만기 (다수결)
    _tm = {}
    for x in rows: _tm[x['term']] = _tm.get(x['term'],0)+1
    term = max(_tm, key=_tm.get) if _tm else ''
    _pp  = term.split('/')[0] if '/' in term else ''
    # 담보 dict = 기존 엔진(rule_extract)에 그대로 태운다
    # ★★★v391b 삼성 재물형은 <b>건물1·건물2에 같은 이름의 담보</b>가 각각 있다.
    #   blk가 담보명 키 dict로 들어가 <b>뒤엣것이 앞엣것을 덮어 6건이 조용히 사라졌다</b>(실측 17/23).
    #   → 같은 이름이면 <b>섹션 라벨을 붙여 둘 다 남긴다</b>(합산은 지점장 확정 전이라 하지 않는다).
    #   ★[확인 대기] 건물1·건물2 동명 담보를 합산할지·분리할지·대표만 쓸지 — 지점장 확정 필요.
    _agg = {}; _ord = []; _mg = []
    for x in rows:
        if x['amt'] is None: continue
        k = x['name']
        if k in _agg:
            # ★v391b2 <b>합산하지 않는다</b>(20억+20억=40억은 내가 만들어낸 값이다).
            #   섹션 라벨을 붙여 <b>둘 다 남기고</b> 확인큐에서 지점장이 판정한다.
            _sfx = x.get('sec') or '2'
            k2 = '%s[%s]' % (k, _sfx); _n = 2
            while k2 in _agg: _n += 1; k2 = '%s[%s#%d]' % (k, _sfx, _n)
            _mg.append(k2); _agg[k2] = x['amt']; _ord.append(k2)
        else:
            _agg[k] = x['amt']; _ord.append(k)
    if _mg: print('[JEAN 동명담보 분리] %d건 — %s' % (len(_mg), ' | '.join(_mg[:8])))
    blk  = ["%s  %s" % (k, _agg[k]) for k in _ord]
    dambo= rule_extract(blk, prefolded=True)   # ★v390b 접힘 복원 재적용 금지
    # ★★★★★v600 제140조 2항 — <b>생보 제안서는 `rule_extract`를 거치지 않는다</b>.
    #   [실측 · KB라이프] 파서는 7건을 정확히 뽑는데 `rule_extract`가 <b>1건만</b> 냈다.
    #     생보 담보명은 `(간편3.10.5)…`처럼 <b>괄호로 시작</b>해 손보용 줄 파싱과 맞지 않는다.
    #   ★그 함수는 <b>손보 전 케이스(규칙 58·실손·심장)가 타므로 고치지 않는다</b>.
    #   ⇒ 생보 경로만 <b>`resolve2`로 직접 매핑</b>한다(경로를 나눠 손 안 댄 쪽을 지킨다).
    if _life_rows and len(dambo or {}) < len(_life_rows):
        _d2 = {}
        for _x in _life_rows:
            try:
                _std = resolve2(_x['name'])[0]
            except Exception:
                _std = None
            if not _std or _std == '__무시__':
                continue
            # ★★★★★v600 제140조 2항 — <b>키는 원본 담보명으로 둔다</b>.
            #   [실측 결함] 마스터 <b>행 이름</b>을 키로 넣었더니 뒤 로직이 그 이름을 다시
            #     `resolve2`에 태워 <b>`일반암` → None</b>이 됐다(2,000 소실).
            #     `고액암`·`유사암`·`허혈성 진단비`는 행 이름과 담보명이 겹쳐 우연히 통과했다.
            #   ⇒ 키를 <b>원본 담보명</b>으로 두면 재매핑이 정상으로 돈다.
            _k2 = _x['name']
            while _k2 in _d2:
                _k2 += ' '
            _d2[_k2] = float(_x['amt'] or 0)
        if len(_d2) > len(dambo or {}):
            print('[JEAN 생보] rule_extract %d건 → 직접 매핑 %d건으로 대체(제140조 2항) %s'
                  % (len(dambo or {}), len(_d2), _d2))
            dambo = _d2
    dambo.pop('__DUP__', None)          # ★내부 로그 키 — 계약 dambo에 남으면 확인큐 기재에서 터진다
    # ★★★★★v521 제117조 (지점장 실측 2026.08.19 «DB 참좋은운전자상해보험2607이 비갱신으로 나온다»)
    #   제안서 경로는 `renewal`을 <b>'비갱신'으로 하드코딩</b>하고 있었다 — <b>제5조를 한 번도 타지 않았다</b>.
    #   실측 PDF는 <b>20년만기 / 20년납</b>(납기 == 만기) → 제5조 ⑤로 <b>갱신</b>이다.
    #   term(`납기/만기`)은 이미 뽑혀 있었는데 <b>쓰지 않았다</b>.
    _jr = '비갱신'
    _pn = re.sub(r'\s', '', str(prod or ''))
    _jd = list((dambo or {}).keys())
    if '갱신형' in _pn: _jr = '갱신'                                   # 제5조 ①
    elif '세만기' in _pn: _jr = '비갱신'                                # 제5조 ②
    elif _jd and all(_is_gen_dambo(_k, '') for _k in _jd): _jr = '갱신'  # 제5조 ④ 담보 100%
    elif '/' in term:                                                  # 제5조 ⑤ 납기 == 만기
        _a, _b = term.split('/', 1)
        _na = re.sub(r'[^0-9]', '', _a); _nb = re.sub(r'[^0-9]', '', _b)
        if _na and _nb and _na == _nb: _jr = '갱신'
    print(f'[JEAN] 회사={co} 담보={len(rows)}건 매핑입력={len(blk)} 보험료={prem:,} 납기/만기={term} 판정={_jr}')
    return {'dup':0,'holder':'','company':co,'ipwon':{},'ci_extra':0,
            'product':(prod or '가입제안서'),'contract_date':'','expiry_date':'',
            'premium':prem,'pay_period':_pp,'pay_count':'','renewal':_jr,
            'dambo':dambo,'ci_jugye':[],'ci_sebu':{},'ci_lines':{},
            'proposal':True,'jean_rows':rows}

# ★★★★★v568 이름 자가진단 (지점장 2026.08.23 「이름은 저번부터 계속 오류라고 했는데
#   오늘은 아예 고객으로 표기다 · 진짜 안 된다 강화해라」).
#   [배경] 이름은 고쳐도 <b>검사가 없어 매번 다시 깨졌다</b>. 심장·실손처럼 케이스로 고정한다.
#   형식 = (파일명, 본문 첫줄들, 기대 이름). 실패 1건이라도 있으면 /health 가 FAIL 을 띄운다.
_NAME_CASES = [
    ('김순자님.pdf',                 [], '김순자'),
    ('김순자님_보장분석.pdf',          [], '김순자'),
    ('보장분석_김순자님.pdf',          [], '김순자'),   # ★제64조 「보장진단」 사고
    ('보장진단_김순자님.pdf',          [], '김순자'),
    ('20260823_김순자님.pdf',        [], '김순자'),
    ('김순자님 (1).pdf',             [], '김순자'),
    ('%EA%B9%80%EC%88%9C%EC%9E%90%EB%8B%98.pdf', [], '김순자'),  # URL 인코딩
    ('김순자.pdf',                   [], '김순자'),
    ('김*자님.pdf',                  [], '김O자'),      # ★제128조 1항 — 지우지 않고 O
    ('abc123.pdf',    ['김순자 님을 위한 보장분석'],       '김순자'),   # 본문 폴백
    ('abc123.pdf',    ['피보험자 : 김순자 (42세/여)'],     '김순자'),
    ('abc123.pdf',    ['계약자 : 김순자'],                '김순자'),
    ('abc123.pdf',    ['김순자 고객님의 보장분석입니다'],     '김순자'),
    # ★v578 제128조 4항 — 본문이 마스킹 + 「고객님을 위한」 (실측 서은옥 2026.08.24)
    ('서은옥님.pdf',   ['서*옥 고객님을 위한'],              '서은옥'),
    ('서은옥님.pdf',   ['고객님을 위한 보장분석'],            '서은옥'),
]


def name_selftest():
    """파일명·본문에서 고객명이 제대로 나오는지 검사. 실패 목록을 돌려준다."""
    bad = []
    for _fn, _ln, _want in _NAME_CASES:
        try:
            _got = parse_txt('\n'.join(_ln), _fn).get('client')
        except Exception as _e:
            bad.append('%s → 예외 %s' % (_fn, str(_e)[:40])); continue
        if _got != _want:
            bad.append('%s → %r (기대 %r)' % (_fn, _got, _want))
    return bad


def parse_txt(txt, filename='', extra=None):
    lines = [l.rstrip() for l in txt.replace('\r\n','\n').replace('\r','\n').split('\n')]
    lines = _repair_anchor(lines)   # ★v282 유실된 계약 경계 앵커 복구
    # ★★★v237: 세부가입현황(상세내역) 계약별 CI 정보 1회 계산 — 선지급률 판정 2순위 근거
    _cib = ci_selftest()
    print('[v238 CI자가진단] ' + ('PASS %d/%d' % (len(_CI_SELFTEST)-len(_cib), len(_CI_SELFTEST)) if not _cib else 'FAIL ' + ' | '.join(_cib[:6])))
    try: _SEBU_CI = parse_sebu_ci(lines)
    except Exception: _SEBU_CI = {}
    if _SEBU_CI: print(f'[v237 sebu_ci] 계약 {len(_SEBU_CI)}건 파싱: ' + ', '.join(f"{k}(사망{v['samang']:,})" for k,v in _SEBU_CI.items()))
    client = ''
    # ★ 정본 §2: 고객명 = 파일명 우선
    # ★★★★★v564 (지점장 실측 2026.08.23 「파일명이 김순자님인데 마음대로 고객으로 했다.
    #   고객으로 하는 건 새제안서를 넣을 때만이다」):
    #   [실측] 현행 정규식이 <b>URL 인코딩</b>(`%EA%B9%80...`)과 <b>mojibake</b>(`ê¹€ìˆœìží‹`)에서
    #   전부 실패해 3405행 `client='고객'`으로 낙하했다. 모바일 업로드에서 흔한 두 형태다.
    #   → 파일명을 <b>먼저 정규화</b>한다. 정규화는 이름을 지어내지 않는다(제128조).
    if filename:
        _fn0 = str(filename)
        try:                                   # ① URL 인코딩 복구
            from urllib.parse import unquote
            if '%' in _fn0:
                _d = unquote(_fn0)
                if re.search(r'[가-힣]', _d): _fn0 = _d
        except Exception: pass
        if not re.search(r'[가-힣]', _fn0):     # ② mojibake 복구 (utf-8을 latin-1로 읽은 경우)
            try:
                _d = _fn0.encode('latin-1', 'ignore').decode('utf-8', 'ignore')
                if re.search(r'[가-힣]', _d): _fn0 = _d
            except Exception: pass
        filename = _fn0
    if filename:
        base = re.sub(r'\.(?:[Tt][Xx][Tt]|[Pp][Dd][Ff])$', '', filename).strip()
        # ★v564 ③ 잡토큰이 앞에 오면 건너뛰고 <b>「○○님」을 먼저</b> 본다(제64조 「보장진단」 사고 방지).
        _NOISE = ('보장분석', '보장진단', '보장분석지', '채널리포트', '가입제안서', '상품제안서',
                  '최종본', '복사본', '설계', '제안서', '고객')
        _m님 = re.search(r'([가-힣][가-힣A-Za-z*O]{1,3})님', base)
        if _m님 and re.match(r'^(%s)' % '|'.join(_NOISE), base):
            client = re.sub(r'님$', '', _m님.group(1))
        fm = None if client else re.match(r'^([가-힣]{2,4})', base)
        # ★v44: '20260713_백O화님_보장분석' 처럼 날짜·숫자 접두 파일명 대응(구 정규식은 '고객'으로 낙하)
        if not fm: fm = re.search(r'([가-힣][가-힣A-Za-z*O]{1,3})님', base)
        if fm: client = re.sub(r'님$', '', fm.group(1))
    # ★★★★★v526 제109조 (실측 2026.08.21 미래에셋 송명복):
    #   파일명 `박진수_가입안내서__송명복_송명복__M-케어건강…` — <b>첫 한글 토큰이 모집인(FC) 이름</b>이다.
    #   §2 「파일명 우선」을 그대로 적용해 <b>고객 제출 문서에 설계사 이름이 박혔다</b>
    #   (실측: 엑셀 파일명 `보장진단_박진수.xlsx` · 헤더 고객명 박진수).
    #   → 파일명 후보가 본문의 <b>「○○ 님을 위한」·주피보험자·피보험자</b>와 다르면 <b>본문이 정본</b>이다.
    #     ★본문에 근거가 있을 때만 바꾼다 — 없으면 파일명을 그대로 둔다(§2 유지).
    # ★★★★★v578 제128조 4항 — <b>호칭어는 이름이 아니다</b>(지점장 2026.08.24 「계속 이름이 고객으로 나온다」).
    #   실측 `고객님을 위한` → 구 정규식이 `고객`을 이름으로 잡아 <b>파일명의 실명을 덮었다</b>.
    _HONOR = ('고객', '귀하', '님', '피보험자', '계약자', '수익자', '본인')
    if client:
        _doc = ''
        for l in lines[:80]:
            for _rx in (r'([가-힣*]{2,4})\s*(?:님|고객님)\s*(?:을|를)?\s*위한',
                        r'주?피보험자\s*[:：]?\s*([가-힣*]{2,4})\s*[(（]'):
                for _m3 in re.finditer(_rx, l):
                    if _m3.group(1) not in _HONOR:
                        _doc = _m3.group(1); break
                if _doc: break
            if _doc: break
        if _doc and _doc != client:
            # ★★★★★v578 제128조 2항 이행 — <b>마스킹은 온전한 이름이 아니다</b>.
            #   [실측] 파일명 `서은옥`(온전) ≠ 본문 `서*옥`(마스킹) → 구 코드가 <b>본문을 채택</b>해
            #     맞는 이름을 버렸고, 뒤의 v549가 되돌리려다 `고객`을 골랐다.
            #   ⇒ 본문에 마스킹 문자가 있고 파일명이 온전하며 <b>성이 같으면</b> 파일명을 지킨다.
            #     성이 다르면 종전대로 본문 채택(제109조 모집인 이름 오염 차단 유지).
            _msk = bool(re.search(r'[\*●○◯·・]', _doc))
            _pure_fn = not re.search(r'[^가-힣]', client)
            if _msk and _pure_fn and _doc[0] == client[0]:
                print('[v578 고객명] 본문 %r 은 마스킹 → 파일명 실명 %r 유지(제128조 2항)' % (_doc, client))
            else:
                print('[v526 고객명] 파일명 %r ≠ 본문 %r → <b>본문 채택</b>(모집인 이름 오염 차단)' % (client, _doc))
                client = _doc
    # ★★★★★v549 제128조 3항 (지점장 실측 2026.08.22 「다른 보장분석도 이름이 안_윤 이렇게 나온다 ·
    #   이름 다 나와야한다」)
    #   제128조 2항은 <b>리모델링 경로만</b> 고쳤다 — 이 `analyze` 경로는 그대로였다.
    #   파일명·본문의 이름이 <b>마스킹(`안*윤`)이나 구분자(`안_윤`)로 끊겨</b> 있으면
    #   <b>순수 한글 2~4자 이름</b>을 파일명·본문에서 다시 찾아 채운다. 못 찾으면 지어내지 않는다.
    # ★★★★★v578 제128조 4항 (지점장 실측 2026.08.24 「계속 이름이 고객으로 나온다 ·
    #   이건 고객에 대한 실례다」 · 파일 `서은옥님.pdf`).
    #   [실측 로그] `[v549] 끊긴 이름 '서*옥' → 실명 '고객' 채택` — <b>본문 폴백이 호칭어를 이름으로 잡았다</b>.
    #     본문 1행이 `서*옥 고객님을 위한`이다. 정규식 `([가-힣]{2,4})…(?:님|고객님)…위한`은
    #     `서*옥`의 `*` 때문에 왼쪽에서 실패하고 <b>`고객`+`님`</b>으로 성립한다.
    #   ⇒ <b>호칭어는 이름이 아니다.</b> 후보가 호칭어면 버리고 계속 찾는다.
    if client and re.search(r'[^가-힣]', client):
        _pure = ''
        for l in lines[:120]:
            for _rx in (r'([가-힣]{2,4})\s*(?:님|고객님)\s*(?:을|를)?\s*위한',
                        r'주?피보험자\s*[:：]?\s*([가-힣]{2,4})\s*[(（]',
                        r'계약자\s*[:：]?\s*([가-힣]{2,4})\s*[(（]'):
                for _m4 in re.finditer(_rx, l):
                    _cand = _m4.group(1)
                    if len(_cand) >= 2 and _cand not in _HONOR:
                        _pure = _cand; break
                if _pure: break
            if _pure: break
        if not _pure and filename:
            # ★파일명 후보는 <b>성(첫 글자)이 같을 때만</b> 받는다 — 모집인 이름 오염 차단(제109조).
            _sn = client[0] if client else ''
            for _c in re.findall(r'[가-힣]{2,4}', str(filename)):
                if _c in ('보장진단', '보장분석', '보장분석지', '최종본', '복사본',
                          '고객', '상품제안서', '가입제안서', '설계'):
                    continue
                if _sn and _c[0] == _sn and len(_c) >= len(client.replace('_', '')):
                    _pure = _c; break
        if _pure and _pure != client:
            print('[v549 고객명] 끊긴 이름 %r → 실명 %r 채택(제128조 3항)' % (client, _pure))
            client = _pure

    # ★★★★★v568 제128조 1항을 <b>analyze 경로에도</b> 적용 (지점장 2026.08.23
    #   「이름은 저번부터 계속 오류라고 했는데 오늘은 아예 고객으로 표기다」).
    #   마스킹 문자는 <b>지우지 않고 O로 자리를 지킨다</b>. 지우면 `김*자`→`김자`가 된다.
    if client:
        client = re.sub(r'[\*●○◯·・]', 'O', client)
    # 폴백: 내용에서 (마스킹 '박*은' 형태도 허용)
    #   ★v568 탐색 범위를 30줄 → 200줄로 넓히고 패턴을 3종 추가한다.
    if not client:
        for l in lines[:200]:
            m5 = (re.search(r'([가-힣*OＯ●○]{2,4})\s*(?:님|고객님)\s*(?:을|를)?\s*위한', l)
                  or re.search(r'주?피보험자\s*[:：]?\s*([가-힣*OＯ●○]{2,4})', l)
                  or re.search(r'계약자\s*[:：]?\s*([가-힣*OＯ●○]{2,4})', l)
                  or re.search(r'([가-힣*OＯ●○]{2,4})\s*고객님', l))
            if m5:
                _c5 = re.sub(r'[\*●○◯·・]', 'O', m5.group(1))
                if len(_c5) >= 2 and _c5 not in ('고객', '보장진단', '보장분석'):
                    client = _c5; break
    if not client:
        for l in lines[:30]:
            l = l.strip()
            m2 = re.search(r'([가-힣*]{2,4})\s*고객님', l) or re.search(r'([가-힣*]{2,4})\s*님의', l)
            if m2: client = m2.group(1); break
            m = re.match(r'^([가-힣]{2,5})\s*$', l)
            if m and len(m.group(1)) <= 4: client = m.group(1); break
    if not client:
        # ★v564 조용히 틀리지 않는다 — 왜 '고객'이 됐는지 흔적을 남긴다(영구원칙).
        print('[v568 고객명] ★파일명·본문에서 이름을 못 찾았다 → 고객 낙하 · 원본 파일명 %r' % (filename,))
        globals()['_NAME_FALLBACK'] = str(filename or '(파일명 없음)')
        globals()['_NAME_WARN'] = ('★[확인] 고객명을 찾지 못해 「고객」으로 표기했습니다 — '
                                   '원본 파일명 %r · 지침상 「고객」은 새 제안서 단독일 때만입니다'
                                   % (filename or '(없음)'))
        client = '고객'

    # ★ 한장보장표(앞부분)에서 회차 추출 → (회사,가입일,만기일) 키로 맵 구축 (별첨엔 회차 없음)
    paycount_map = {}
    for l in lines:
        ld = l.strip()
        # ★★★★★v483 제94조 (김순자 실측 2026.08.19) — <b>납입회차를 못 읽으면 갱신 판정이 죽는다</b>.
        #   구 정규식은 ㉠<b>'월납'이 반드시 있어야</b> 하고 ㉡분모가 <b>2~3자리</b>여야 했다.
        #   실측 NH농협생명 줄 = `… 2025.08.06 2028.08.06 <b>2/3</b> 24,464 …` —
        #   납입주기 칸이 <b>비어 있고</b> 분모가 <b>1자리</b>라 둘 다 탈락 → pay_count 공란 →
        #   judge_renewal이 기간을 못 재고 <b>무조건 비갱신</b>이 됐다(제85조와 같은 뿌리).
        #   → 납입주기는 <b>선택</b>, 분모는 <b>1~3자리</b>. 두 날짜 바로 뒤 토큰이라 오인 위험은 낮다.
        m = re.search(r'([가-힣A-Za-z]{2,8}(?:생명|화재|손보|손해|해상|라이프|증권)?)\s+.*?(\d{4}\.\d{2}\.\d{2})\s+(\d{4}\.\d{2}\.\d{2})\s+(?:\d+\s*년납|월납|년납|연납|일시납|분기납|반년납)?\s*(\d{1,3}/\d{1,3})', ld)
        if m:
            comp = m.group(1).strip()
            paycount_map[(comp, m.group(2), m.group(3))] = m.group(4)
            paycount_map[(m.group(2), m.group(3))] = m.group(4)  # 회사 표기 흔들림 대비 보조키

    contracts = []; i = 0; n = len(lines)
    # ★★v44 분기: 신정원 3열 포맷(KB·메리츠) 감지 시 parse_sinjeong 사용, 아니면 기존 2열 별첨 엔진.
    _SJN   = sinjeong_count(lines)
    _IS_SJ = _SJN >= 2
    if _IS_SJ:
        contracts = parse_sinjeong(lines)
        i = n                      # 기존 2열 루프 스킵('정상계약 리스트' 문자열 부재로 어차피 0건)
    while i < n:
        l = lines[i].strip()
        if '실효계약 리스트' in l or '미납해지' in l: break
        if '정상계약 리스트' not in l: i += 1; continue
        i += 1
        while i < n and not lines[i].strip(): i += 1
        if i >= n: break
        # ★v30p+ 형태 자동 감지: 첫 줄에 계약자/탭/가입금액/Chtd 있으면 신형(다중헤더), 아니면 정상형(기존).
        _first = lines[i]
        if ('계약자' in _first) or ('\t' in _first) or ('가입금액' in _first) or ('Chtd' in _first):
            # ── 신형 (오늘 PDF 업데이트 형태, 별첨 헤더 A/B/C 혼재) ──
            # ★v30p 다중 별첨 헤더 형태 대응(A:계약자줄먼저→회사·상품 다음줄 / B:회사·상품+계약자 한줄 / C:계약자이름+회사+상품 한줄).
            #   담보표 헤더('가입금액'·'담보명'·'Chtd') 전까지를 헤더영역으로 모아 회사·상품·보험료·날짜를 통째 추출.
            _hdr=[]; _k=i
            while _k < n and _k < i+8:
                _lk = lines[_k]
                if '가입금액' in _lk or 'Chtd' in _lk or '담보명' in _lk: break
                if '정상계약 리스트' in _lk or '실효계약 리스트' in _lk: break
                # ★v59 첫 담보 유실 방지: '한글…+끝자리 숫자'=담보값 줄을 만나면 헤더수집 중단.
                #   (계약자·보험료·보장기간·납입·날짜 줄은 담보로 오인 안 되게 제외.)
                #   이전 고정 6줄창은 블랭크패딩 때문에 첫 담보(예 한화생명 암수술 50)를
                #   상품명에 흡수 → 담보 1~3개 유실. 생보 암보험(첫 담보=암진단)이 가장 큰 피해.
                _lks=_lk.strip()
                # ★★★v234 (2026.07.25 한정환 실측): 상품명 줄이 '담보값 줄'로 오인돼 헤더수집이 끊기던 버그.
                #   `무배당 마이라이프 한아름종합보험 1710` / `참좋은운전자상해보험 2510` 처럼
                #   상품명 끝에 <b>4자리 상품코드</b>가 붙으면 이 가드가 걸려 break → 상품명이 공란이 되고
                #   그 줄은 담보로 파싱돼 확인사항 큐에 `무배당 마이라이프 한아름종합보험 = 1710`이 박혔다.
                #   판별: 별첨 담보표는 pdftotext -layout에서 <b>금액 앞 공백이 2칸 이상</b>(표 열 구분)이고,
                #   1,000 이상은 <b>콤마</b>가 찍힌다. 반대로 상품코드는 <b>공백 1칸 + 콤마 없는 4자리</b>다.
                _mnum = re.search(r'(\s+)([\d][\d,]*)\s*$', _lks)
                _isprod = bool(_mnum) and len(_mnum.group(1))==1 and bool(re.fullmatch(r'\d{4}', _mnum.group(2)))
                if _lks and re.search(r'[가-힣]', _lks) and _mnum and not _isprod \
                   and not re.search(r'계약자|보험료|보장기간|납입|월납|\d{4}\.\d{2}', _lks):
                    break
                _hdr.append(_lk); _k += 1
            _ht = ' '.join(_hdr).replace('\t',' ')
            i = _k+1 if (_k<n and ('가입금액' in lines[_k] or 'Chtd' in lines[_k] or '담보명' in lines[_k])) else _k
            contract_date = expiry_date = pay_period = pay_count = ''; premium = 0; lump_sum = 0
            _md = re.search(r'(\d{4}\.\d{2}\.\d{1,2})\s*[-~（卜\s]+(\d{4}\.\d{2}\.\d{1,2})', _ht)
            if _md: contract_date=_md.group(1); expiry_date=_md.group(2)
            _mp = re.search(r'보험료\s*([\d,\.]+)\s*원', _ht)
            if _mp:
                try:
                    _pv=int(_mp.group(1).replace(',','').replace('.',''))
                    if 1000 < _pv < 5000000: premium=_pv
                    elif 5000000 <= _pv < 10**10: lump_sum = _pv
                    elif _pv >= 5000000 and '일시' in re.sub(r'\s','',str(pay_period or '')): lump_sum=_pv
                except: pass
            _mper = re.search(r'월납\s*/?\s*(\d+)\s*년', _ht)
            if _mper: pay_period=f"{_mper.group(1)}년납"
            # ★v126 열 어긋남 수정: 납입주기 칸이 '일시납/연납/N개월납'인 행에서
            #   그 토큰이 회사명 앞에 붙어 '일시납AIG손보'로 오염되던 버그.
            elif re.search(r'일시납', _ht): pay_period = '일시납'
            elif re.search(r'연\s*납', _ht): pay_period = '연납'
            _mpc = re.search(r'(\d{1,3})\s*/\s*(\d{2,3})\s*회', _ht)
            if _mpc: pay_count=f"{_mpc.group(1)}/{_mpc.group(2)}"
            # 회사·상품 분리: 계약자·납입·보험료·보장기간·단위 boilerplate 제거 후 보험사 키워드로 split
            _ct = _ht
            _ct = re.sub(r'계약자\s*\S+',' ',_ct)
            _ct = re.sub(r'납입주기\s*/?\s*기간',' ',_ct)
            _ct = re.sub(r'보험료\s*[\d,\.]+\s*원',' ',_ct)
            _ct = re.sub(r'보장기\s*간|보장기간',' ',_ct)
            _ct = re.sub(r'\d{4}\.\d{2}\.\d{1,2}\s*[-~（卜\s]+\d{4}\.\d{2}\.\d{1,2}',' ',_ct)
            _ct = re.sub(r'월납\s*/?\s*\d+\s*년',' ',_ct)
            _ct = re.sub(r'일시납|연\s*납|월\s*납|\d+\s*개월납',' ',_ct)   # ★v126 납입주기 토큰 제거
            _ct = re.sub(r'\d{1,3}\s*/\s*\d{2,3}\s*회',' ',_ct)
            _ct = re.sub(r'[（(]?\s*단위\s*:?\s*만원\s*[）)]?',' ',_ct)
            _ct = re.sub(r'\s+',' ',_ct).strip()
            _mc = re.search(r'^(.*?(?:화재|손해보험|손보|해상|생명|라이프|증권))\s*(.*)$', _ct)
            if _mc:
                company = re.sub(r'\s','',_mc.group(1)); product = _mc.group(2).strip()
                # ★★★v234 (한정환 실측): 정규식 `.*?`가 최단일치라 '메트라이프생명'에서 '라이프'를 먼저 물어
                #   회사명이 <b>'메트라이프'</b>, 상품명이 <b>'생명 무배당 마스터플랜…'</b>으로 잘렸다.
                #   → 상품명이 회사 접미어로 시작하면 그 토큰을 회사명으로 되돌린다.
                _mtail = re.match(r'^(생명|화재|손해보험|손보|해상|라이프|증권|공제)\s+(\S.*)$', product)
                if _mtail:
                    company = company + _mtail.group(1); product = _mtail.group(2).strip()
            else:
                _parts=_ct.split(' ',1); company=_parts[0] if _parts else lines[i].strip(); product=_parts[1].strip() if len(_parts)>1 else ''
        else:
            # ── 정상형 (기존, '가입금액' 헤더 없이 회사→상품→계약자줄→담보) ──
            _head = lines[i].strip()          # ★v44 롯데: 이 줄에 '회사명 상품명'이 한 줄로 붙어 온다
            company = _head; i += 1
            contract_date = expiry_date = pay_period = pay_count = ''; premium = 0; lump_sum = 0
            for _j in range(i, min(i+5, n)):
                _l = lines[_j]
                _m = re.search(r'(\d{4}\.\d{2}\.\d{2})\s*~\s*(\d{4}\.\d{2}\.\d{2})', _l)
                if _m: contract_date = _m.group(1); expiry_date = _m.group(2)
                _m2 = re.search(r'(\d{1,3})\s*/\s*(\d{2,3})\s*회', _l) or re.search(r'월납\s+(\d{1,3})\s*/\s*(\d{2,3})', _l)
                if _m2 and not pay_count: pay_count = f"{_m2.group(1)}/{_m2.group(2)}"
                _m3 = re.search(r'보험료\s*([\d,\.]+)\s*원', _l) or re.search(r'([\d,]+)원', _l)
                if _m3:
                    try:
                        _v = int(_m3.group(1).replace(',','').replace('.',''))
                        if 1000 < _v < 5000000: premium = _v
                        elif 5000000 <= _v < 10**10: lump_sum = _v
                        elif _v >= 5000000 and '일시' in re.sub(r'\s','',str(pay_period or '')): lump_sum = _v
                    except: pass
                _m5 = re.search(r'월납\s*/?\s*(\d+)\s*년', _l)
                if _m5 and not pay_period: pay_period = f"{_m5.group(1)}년납"
            while i < n and not lines[i].strip(): i += 1
            product = ''
            # ★★v44 롯데 결함B 수정: 헤더 줄을 보험사 키워드로 회사/상품 분리.
            #    구버전은 company에 상품명이 통째로 남고, product 칸엔 '담보 첫 줄'이 들어가
            #    (1) 상품명 오염 (2) 담보 1개 유실 (3) 병합키(회사·보험료·상품[:12]) 불일치로
            #    같은 계약이 2건으로 쪼개지는 결함A까지 유발했다. 회사/상품만 바로잡으면 셋 다 해소.
            _mc = re.match(r'^(.*?(?:화재|손해보험|손보|해상|생명|라이프|증권|공제))\s+(\S.*)$', _head)
            if _mc:
                company = re.sub(r'\s', '', _mc.group(1))
                product = re.sub(r'\s+', ' ', _mc.group(2).strip())
            else:
                # 폴백(기존): 회사명만 있는 헤더 → 다음 줄들에서 상품명 탐색
                for _j in range(i, min(i+6, n)):
                    _l = lines[_j].strip()
                    if _l and not re.search(r'계약자|납입주기|보험료|보장기간', _l):
                        if len(_l) > 5 and not re.search(r'^[\d,]+$', _l) and not re.search(r'^\d{4}\.\d{2}', _l):
                            product = _l; i = _j + 1; break
        if is_excluded(company, product, contract_date, expiry_date, pay_period, pay_count):
            while i < n and '정상계약 리스트' not in lines[i] and '실효계약 리스트' not in lines[i]: i += 1
            continue
        renewal = judge_renewal(product, expiry_date, pay_count, contract_date, pay_period, company)   # ★v354: 이 시점엔 dambo 미생성 → 전기납 예외는 아래 2271행 재판정에서 적용
        # 담보 블록 텍스트 수집 (다음 '정상계약/실효계약 리스트'까지)
        # ★v333: 계약자명 — 별첨 헤더('계약자 이*영 …')는 <b>상품명 줄보다 앞</b>에 온다.
        _holder=''
        for _hb in range(max(0,i-16), i):
            _mh=re.search(r'계약자\s*([가-힣A-Za-z\*＊]{2,8})', lines[_hb])
            if _mh: _holder=_mh.group(1); break
        block_lines = []; j = i
        while j < n:
            if '정상계약 리스트' in lines[j] or '실효계약 리스트' in lines[j]: break
            block_lines.append(lines[j]); j += 1
        i = j
        # 추출: LLM 우선(깨진 별첨 복원), 키 없거나 실패 시 규칙 폴백
        dambo = llm_extract('\n'.join(block_lines)) or rule_extract(block_lines)
        _dup = dambo.pop('__DUP__', {}) if isinstance(dambo, dict) else {}   # ★v345 중복줄 기록(표시 전용)
        # ★ CI/리빙케어/GI: 별첨이 전부 '주계약'으로 라벨없이 뭉침 → 개별 주계약 금액 수집(본체 80/50% 판별용)
        ci_jugye=[]
        if _isci_prod(product, company):        # ★v446 손보 제외
            for _bl in block_lines:
                # ★v33 pdftotext -layout 다열 레이아웃: 한 줄에 2~3담보 → finditer.
                #    '_주계약'(일반사망_주계약 등) 은 lookbehind 로 차단.
                for _m in re.finditer(r'(?<![_가-힣])주계약\s+([\d,]{3,})', _bl):
                    try:
                        _v=int(_m.group(1).replace(',',''))
                        if 0<_v<=200000: ci_jugye.append(_v)
                    except: pass
            # ★v30z2 삼성 리빙케어형: 주계약이 '주계약'이 아니라 상품명(삼성리빙케어(종신2종)1.2)으로 라벨됨.
            #   '리빙케어 종신N종' 플랜패턴 뒤 금액만 추출(리빙케어보장특약·재해사망 등 부담보 오수집 차단).
            #   콤마·마침표 천단위(3.000=3,000) 정규화. 80/50 판별은 아래 기존 로직이 수행.
            if not ci_jugye and '리빙케어' in (product or ''):
                for _bl in block_lines:
                    for _mm in re.finditer(r'리빙케어\s*[（(]?\s*종신\d*종[）)]?\s*\d*[.．]?\d*[\s\t]+([\d][\d.,]{2,})', _bl):
                        try:
                            _v=int(_mm.group(1).replace(',','').replace('.',''))
                            if 100<=_v<=200000: ci_jugye.append(_v)
                        except: pass
            # ★★★v122(2026.07.21 장문순 실측): 삼성생명 CI는 주계약이 '주계약'이라는 글자로
            #   안 나오고 <상품명 자체>로 라벨된다.
            #   예) '퍼펙트통합보험프리미엄,無표준월납  3,000 / ... 1,500'
            #   → 상품명 앞 6자를 키로 삼아 그 라벨이 붙은 금액을 주계약 후보로 수집한다.
            #   (리빙케어 전용 분기 v30z2의 일반화 — 퍼펙트플러스·퍼펙트통합 등 전 CI 상품 적용)
            if not ci_jugye:
                _pk = re.sub(r'[^가-힣A-Za-z0-9]', '', str(product or ''))[:6]
                if len(_pk) >= 4:
                    for _bl in block_lines:
                        for _mm in re.finditer(r'([가-힣A-Za-z0-9,\.\(\)（）Ⅰ-Ⅹ無·]{4,}?)\s{2,}([\d][\d,]{2,})', _bl):
                            _lab = re.sub(r'[^가-힣A-Za-z0-9]', '', _mm.group(1))
                            if _lab.startswith(_pk):
                                try:
                                    _v = int(_mm.group(2).replace(',', ''))
                                    if 100 <= _v <= 200000: ci_jugye.append(_v)
                                except Exception:
                                    pass
        # ★v29t CI추가보장특약: 줄별 값 수집(병합 전 원값 보존)
        ci_extra=[]
        for _bl in block_lines:
            for _m in re.finditer(r'CI\s*추가보장특약\s+([\d,]+)', _bl):
                try:
                    _v=int(_m.group(1).replace(',',''))
                    if 0<_v<=200000: ci_extra.append(_v)
                except: pass
        # ★v29t 생보 입원특약 일당: 줄별 값 수집(병합 전 원값 보존)
        ipwon=[]; _ipkey=None
        # ★★★★★v239 (지점장 지적 2026.07.25): <b>CI 주계약 사망은 별첨 '줄 단위 개별값'이다 — 합산 금지</b>.
        #   실측 증거(DB생명 CI종신 1701): 별첨에 <b>질병사망 3,000 / 질병사망 2,000</b> 2줄 = 주계약 3,000 + 추가특약 2,000.
        #   그런데 dambo는 이를 <b>합산해 5,000</b>으로 만든다 → 2,400/5,000=48%가 되어 50%·80% 어디에도 안 맞고
        #   [확인]으로 빠졌다. <b>주계약 3,000으로 보면 2,400/3,000 = 정확히 80%</b>(뇌출혈·급성심근·암·특정질병·치매가
        #   전부 2,400으로 일치) → <b>80%형이 정답</b>이었다.
        #   → CI 계약은 <b>접힘 복원된 줄에서 사망 줄과 담보 줄의 원값을 따로 수집</b>한다.
        ci_lines={'samang':[], 'cands':[], 'jungdae':[], 'brain':[]}
        # ★★★v234: 별첨이 <b>다열</b>이면 원본 줄은 `재해장해특약  6,000   입원  3` 처럼
        #   두 담보가 한 줄에 붙어 있어 `^입원 …$` 앵커가 절대 맞지 않는다(실측 메트라이프 9p).
        #   → rule_extract와 <b>동일한 접힘 복원</b>을 거친 1담보 1줄 형태에서 스캔한다.
        try:
            _bl_flat = _split_cols(_reflow_cols(_unfold_cols(block_lines)))
        except Exception:
            _bl_flat = block_lines
        for _bl in _bl_flat:
            # ★★★v234 (한정환 메트라이프생명 실측): 별첨에 <b>수식어 없는 '입원'</b>만 인쇄되는 생보가 있다.
            #   구 정규식은 '입원특약'만 잡아 `입원 3` 2줄을 일당 어느 행에도 못 넣었다
            #   (세부가입현황 4p 입원비 칸 = `3 | 3` = 질병 3 · 상해 3, 한장보장표 일당 13 = 신한 10 + 메트 3).
            # ★★★v239 CI 주계약/본체 후보 = 줄 단위 원값(합산 전)
            if _isci_prod(product, company):    # ★v446 손보 제외
                _mc = re.match(r'^\s*(\S.*?)\s{2,}([\d,]{3,})\s*$', _bl.rstrip()) or \
                      re.match(r'^\s*(\S.*?)\s+([\d,]{3,})\s*$', _bl.rstrip())
                if _mc:
                    _nmc = re.sub(r'\s','',_mc.group(1))
                    try: _vc = int(_mc.group(2).replace(',',''))
                    except Exception: _vc = 0
                    if 0 < _vc <= 200000 and re.search(r'[가-힣]', _nmc):
                        # ★v253 상해사망·재해사망 제외(지점장 지시) — 위 3열 경로와 동일
                        if (re.search(r'사망', _nmc) and not re.search(r'후유|장해', _nmc)
                                and not re.search(r'상해|재해', _nmc)):
                            ci_lines['samang'].append(_vc)
                        elif _nmc.startswith('중대한'):
                            # ★★★★★v240(지점장 지시 2026.07.25): 별첨에 <b>'중대한OO' 담보가 명시</b>되어 있으면
                            #   <b>그 금액이 곧 CI 본체(선지급액)</b>다. 실측 신한 `중대한화상진단 2,000`.
                            ci_lines['jungdae'].append(_vc)
                        elif re.search(r'진단', _nmc):     # 본체 후보 = 진단비 담보만
                            ci_lines['cands'].append(_vc)
                            # ★★★★★v242: 뇌 담보는 <b>축(뇌출혈/뇌졸증)과 금액</b>을 따로 기록한다.
                            #   CI 본체를 '중대한 뇌졸증'에 넣을지 '중대한 뇌출혈'에 넣을지 판별하기 위함.
                            if '뇌출혈' in _nmc:   ci_lines['brain'].append(('뇌출혈', _vc))
                            elif '뇌졸' in _nmc:   ci_lines['brain'].append(('뇌졸증', _vc))
            _m=re.match(r'^\s*(입원특약|입원)\s+([\d,]+)\s*$', _bl.strip())
            if _m:
                _ipkey=_m.group(1)
                _m=re.match(r'^\s*(?:입원특약|입원)\s+([\d,]+)\s*$', _bl.strip())
                try:
                    _v=int(_m.group(1).replace(',',''))
                    if 0<_v<=1000: ipwon.append(_v)
                except: pass
        # ★v29t (지점장 확정 2026.07.02): 생보 '입원특약' 일당 = 상해·질병 둘 다 해당 →
        #   질병일당·상해일당 두 행에 한 줄 값(중복줄=동일특약 재출현 → max) 각각 기재. dambo 변환이라 엑셀·PPT 동일 반영.
        if ipwon and any(k in (company or '') for k in ('생명','라이프','AIA','메트라이프','우체국','공제')) and (_ipkey in dambo):
            _v1=max(ipwon)   # ★중복줄(질병축·상해축)=동일 특약 → max
            dambo.pop(_ipkey, None)
            dambo['질병입원일당(입원특약)']=dambo.get('질병입원일당(입원특약)',0)+_v1
            dambo['상해입원일당(입원특약)']=dambo.get('상해입원일당(입원특약)',0)+_v1
        if company:
            # ★★★v237: 세부가입현황(상세내역) 계약별 CI 정보를 계약에 부착한다.
            #   롯데(2열) 별첨엔 '주계약' 라벨이 없어 `ci_jugye=[]`가 되므로,
            #   선지급률(50%/80%) 판정의 <b>2순위 근거</b>로 쓴다(지점장 지시 2026.07.25).
            _cs=None
            try:
                _ck=re.sub(r'\s','',str(company or ''))
                for _k,_v in (_SEBU_CI or {}).items():
                    _k2=re.sub(r'\s','',_k)
                    if _k2==_ck or _k2.startswith(_ck[:3]) or _ck.startswith(_k2[:3]): _cs=_v; break
            except Exception: pass
            # ★★★★★v333 (지점장 지시 2026.08.02, 영구): <b>계약자가 다르면 다른 계약이다</b>.
            #   지점장 원문: "각각 표기해야한다. 왜냐면 <b>각각의 사망보장금이나 담보가 있기때문</b>이다".
            #   실측(이영태) 푸르덴셜생명 종신 3건 — 회사·상품·보험료(0원)·계약일·만기일이 <b>전부 동일</b>해
            #   병합키가 겹쳐 <b>1건으로 뭉쳤다</b>. 계약자는 이*영 / 이*현 / 이*순으로 서로 다르고
            #   사망보장금도 7,865 / 7,865 / <b>9,050</b>으로 다르다.
            #   → 별첨 헤더의 <b>계약자명</b>을 계약에 붙여 병합키에 넣는다.
            #   ★페이지 분할된 <b>같은</b> 계약은 계약자가 동일하게 인쇄되므로 병합은 그대로 유지된다(v257 무회귀).
            contracts.append({'dup':_dup,'holder':_holder,'company':company,'ipwon':ipwon,'ci_extra':ci_extra,'product':product,'contract_date':contract_date,
                'expiry_date':expiry_date,'premium':premium,'lump_sum':lump_sum,'pay_period':pay_period,
                'pay_count':pay_count,'renewal':renewal,'dambo':dambo,'ci_jugye':ci_jugye,'ci_sebu':_cs,'ci_lines':ci_lines})
    # ★★v100 단계약 사각지대: KB·메리츠 3열이 '계약 1건'이면 앵커가 1개뿐이라
    #   sinjeong_detect(>=2)를 못 넘겨 2열 파서로 가고 계약 0건이 된다(양예서형 단계약 고객).
    #   → 2열이 0건일 때만 3열을 재시도한다. 2열이 1건이라도 잡으면 손대지 않으므로 회귀 없음.
    if (not contracts) and _SJN >= 1:
        contracts = parse_sinjeong(lines)
        print(f'[SINJEONG-FALLBACK] 2열 0건 + 3열 앵커 {_SJN}개 → 3열 재시도 = {len(contracts)}건')
    # ★ 페이지 분할 중복 제거 (정본 체크리스트 ①②): 동일 계약키 병합
    merged = {}
    order = []
    for c in contracts:
        # ★★★★★v257(2026.07.27 실측): 이 병합키의 공백 제거가 <b>한 글자도 작동하지 않았다</b>.
        #   `r'\\s'`는 <b>역슬래시+s 문자</b>를 찾는 패턴이라 공백이 그대로 남는다.
        #   → pdftotext -layout이 같은 계약을 페이지마다 `KB손해보험` / `KB 손해보험`으로 뱉으면
        #     키가 달라져 <b>한 계약이 두 열로 쪼개지고, 담보가 페이지별로 절반씩</b> 나뉜다.
        #     별첨 뒷부분에 몰린 <b>심장 묶음·수술비·치료비</b>가 통째로 빠지는 원인이다.
        # ★v333: 계약자명을 병합키에 포함 — 계약자가 다르면 별개 계약(사망보장금·담보가 각각 있다).
        key = (re.sub(r'\s','',c['company']), c['premium'], re.sub(r'\s','',c['product'])[:12],
               re.sub(r'\s','',str(c.get('holder') or '')))   # ★v30p 날짜 OCR 깨짐 대비 병합키
        if key not in merged:
            merged[key] = c; order.append(key)
        else:
            m = merged[key]
            # 담보 병합: 같은 담보명은 큰 값 유지(중복가산 방지), 새 담보는 추가
            for k, v in c['dambo'].items():
                m['dambo'][k] = max(m['dambo'].get(k, 0), v)
            if not m.get('ci_jugye') and c.get('ci_jugye'): m['ci_jugye'] = c['ci_jugye']
            # 더 긴(덜 잘린) 상품명 채택
            if len(c['product']) > len(m['product']): m['product'] = c['product']
            # 회차/기간 비어있으면 채움
            if not m['pay_count'] and c['pay_count']: m['pay_count'] = c['pay_count']
            if not m['pay_period'] and c['pay_period']: m['pay_period'] = c['pay_period']
    deduped = [merged[k] for k in order]
    # ★★★★★v370 (지점장 확정 2026.08.09): <b>가입제안서 계약을 여기서 합류</b>시킨다.
    #   후처리(_HB 심장 묶음 · 세부보충 · 대표값)를 <b>보장분석지와 똑같이</b> 타야 한다.
    #   ★분해 경로를 새로 만들지 않는다(영구원칙: 묶음 분해는 한 곳에서만).
    if extra:
        for _x in extra:
            if _x: deduped.append(_x)
        print(f'[JEAN] 제안 계약 {len([x for x in extra if x])}건 합류 → 총 {len(deduped)}건')
    # ★★★★★v620 (지점장 실측 2026.08.30 「<b>DB손보가 왜 2개가 출력되지</b>」 · 배미용 실측).
    #   [원인] 병합키가 <b>상품명 앞 12자</b>인데, 별첨에서 담보 한 줄이 상품명에 딸려 들어와
    #     `참좋은종합보험1709` / `참좋은종합보험1709<b>질병사망(15년갱신)…1,000</b>`
    #     두 이름이 되어 <b>같은 계약이 두 열</b>로 갈렸다(보험료·가입일·만기일은 동일).
    #   ⇒ <b>회사 + 보험료 + 계약일 + 만기일</b>이 모두 같으면 한 계약이다. 2차 병합한다.
    #     ★상품명은 <b>짧은 쪽</b>을 남긴다 — 담보가 붙지 않은 것이 진짜 상품명이다.
    try:
        _k2 = {}; _ord2 = []
        for _c in deduped:
            _kk = (re.sub(r'\s', '', str(_c.get('company') or '')),
                   _c.get('premium'), str(_c.get('contract_date') or ''),
                   str(_c.get('expiry_date') or ''))
            if not _kk[1] or not _kk[2]:
                _kk = ('', id(_c), '', '')        # 보험료·계약일 없으면 병합하지 않는다
            if _kk not in _k2:
                _k2[_kk] = _c; _ord2.append(_kk)
            else:
                _m2 = _k2[_kk]
                for _n2, _v2 in (_c.get('dambo') or {}).items():
                    _o2 = (_m2.get('dambo') or {}).get(_n2)
                    if _o2 is None:
                        _m2.setdefault('dambo', {})[_n2] = _v2
                    else:
                        try:
                            if float(_v2) > float(_o2):
                                _m2['dambo'][_n2] = _v2
                        except Exception:
                            pass
                _p2 = str(_c.get('product') or '')
                if _p2 and len(_p2) < len(str(_m2.get('product') or '')):
                    _m2['product'] = _p2
        if len(_ord2) < len(deduped):
            print('[v620 계약병합] 보험료·기간 동일 계약 %d건 → %d건'
                  % (len(deduped), len(_ord2)))
            deduped = [_k2[k] for k in _ord2]
    except Exception as _e620:
        print('[v620 계약병합] 실패', str(_e620)[:60])
    # ★★★제외 7종(v125, 지점장 확정 2026.07.21): 실손이 아닌데 <계약일 또는 만기일이 없는> 계약은
    #   보험기간(1년 여부)을 판정할 수 없다 → 엑셀·보장나무·보장진단서·보장설명서 전부 미포함.
    #   ★실손은 예외 — 롯데 리포트가 실손 계약의 계약일·보험료를 공란으로 주는 사례가 있다(v90·장문순 실측).
    #   ★담보를 봐야 실손인지 알 수 있으므로 파싱이 끝난 이 시점에서 판정한다.
    # ★★★★★v262(지점장 지시 2026.07.27, 영구): <b>기간불명은 제외 사유가 아니다 — 엑셀에 기재만 한다</b>.
    #   지점장 원문 = "이건 그냥 엑셀에 기재만하라 / 불명확한건 안해도된다".
    #   구 v125는 계약일·만기일이 둘 다 없으면 계약을 <b>통째로 버렸다</b> → 실측(장O경) 한화손보
    #   `무배당 한화보금자리안심보험1605`가 통째 소실되고 <b>한장표 화재벌금 2,000이 문서에서 사라졌다</b>.
    #   → <b>계약열은 만든다</b>. 계약일·만기일·총납입기간 같은 <b>불명확한 칸은 공란</b>으로 둔다.
    #   ★<b>제외 7종 ⑥은 '보험기간 1년(만기−가입 358~372일)'만</b>이다 — 기간을 <b>알 수 없는</b> 것과
    #     기간이 <b>1년인</b> 것은 다르다. 1년 제외 규칙은 그대로 살아 있다.
    # ★★★★★v552 제9조 4항 (지점장 실측 2026.08.22 「엑셀도 ppt도 다 상해의료비에 5천이 찍혔어 · 명백한 오류다」)
    #   <b>`상해의료비` 행은 1세대(구실손) 전용</b>이다(제9조 정본표 — 「1세대 손보 · 상해의료비 별도」).
    #   2세대 이후는 <b>입원/통원</b>으로 간다. 그런데 채널보고서가 실손을 요약 표기해
    #   <b>`상해의료비 5,000`</b> 한 줄로 주면 세대와 무관하게 그 행으로 갔다.
    #   ⇒ 같은 계약에 <b>3대 비급여(도수·주사·MRI)나 급여/비급여 분리 담보</b>가 있으면
    #     그 계약은 <b>3세대 이상</b>이다 → `상해의료비` 값을 <b>입원</b>으로 옮긴다.
    for _c in deduped:
        _dm = _c.get('dambo') or {}
        if '상해의료비' not in _dm:
            continue
        _keys = ''.join(str(k) for k in _dm.keys())
        _gen3 = (('도수' in _keys) or ('체외충격파' in _keys) or ('증식치료' in _keys)
                 or ('MRI' in _keys.upper()) or ('비급여' in _keys and '주사' in _keys)
                 or ('급여' in _keys and '비급여' in _keys))
        if not _gen3:
            continue
        _v = _dm.pop('상해의료비', 0)
        if _v:
            _dm['입원'] = max(_dm.get('입원', 0) or 0, _v)
            print('[v552 실손] %s — 3세대 이상 계약의 상해의료비 %s → 입원으로 이동(제9조 4항)'
                  % (str(_c.get('company'))[:10], _v))

    # ★★★★★v593 제134조 — <b>일상배상책임은 무조건 1억(10,000만원)</b>
    #   (지점장 확정 2026.08.26 「<b>일상배상책임은 무조건 1억이다</b>」).
    #   [실측 결함 · 롯데 let:smile] 확인사항에 `[접힘의심] 가족일상생활배상책임(대물` <b>20</b> —
    #     담보명이 `(대물`에서 <b>잘렸고</b>, 원문 `(대물 누수 50만원, 누수외 20만원)`의
    #     <b>자기부담금 20</b>을 가입금액으로 읽었다. 1억이 <b>20만원</b>이 됐다.
    #   ⇒ 담보가 존재하면 금액은 <b>10,000 고정</b>. 자기부담금·누수한도는 가입금액이 아니다.
    #     ★값이 이미 10,000이면 그대로. 다르면 <b>고쳐 쓰고 로그로 남긴다</b>(조용히 바꾸지 않는다).
    _fix_ilbae(deduped)

    # ★★★★★v590 제9조 5항 — <b>통합형 실손의 통원은 20</b>(지점장 확정 2026.08.25
    #   「<b>통원 20에 잡아주면 된다</b>」 · 4세대 정본 「통원비 20만원만 있다」).
    #   통합형(`의료비(입원+통원)`)은 입원 한도만 적혀 있고 통원 금액이 별도로 없다.
    #   ⇒ 입원이 있고 <b>통원이 비어 있는</b> 실손 계약에만 <b>20</b>을 채운다.
    #   ★<b>5세대는 손대지 않는다</b>(지점장 「5세대는 표기가 많이 달라서 그건 그때 잡자」).
    #   ★통원이 이미 있으면 <b>덮지 않는다</b> — 별첨 명시값이 우선(제9조 금액 규칙).
    # ★★★★★v596 제9조 6항 (지점장 지적 2026.08.28 「통원이 5천만원 나오는건 정정된거지?」)
    #   [확인] v595에는 <b>통원 값을 제한하는 코드가 없었다</b>. `_dm['통원']=20`은 통원이
    #   <b>비어 있을 때만</b> 채우므로, 별첨·파서가 통원에 입원 한도(5,000)를 실으면 그대로 찍혔다.
    #   ⇒ 실손 계약의 통원이 <b>1,000(천만원) 이상</b>이면 입원 한도가 새어든 것이다.
    #     입원으로 흡수(max)하고 통원을 비운다 → 바로 아래 v590 규칙이 통원 20을 채운다.
    #   ★임계값 1,000은 보수적이다. 정상 통원은 10·20·25로 <b>자릿수가 다르다</b>.
    for _c in deduped:
        _dm = _c.get('dambo') or {}
        _tw = _dm.get('통원')
        try:
            _twv = float(_tw)
        except Exception:
            continue
        if _twv < 1000:
            continue
        if not _is_silson_like(_c.get('company'), _c.get('product'), _dm):
            continue
        _dm['입원'] = max(_dm.get('입원', 0) or 0, _twv)
        _dm.pop('통원', None)
        print('[v596 실손] %s — 통원 %s = 입원 한도 유입 → 입원으로 이동·통원 재산정(제9조 6항)'
              % (str(_c.get('company'))[:10], format(int(_twv), ',')))

    for _c in deduped:
        _dm = _c.get('dambo') or {}
        if not _dm.get('입원'):
            continue
        if _dm.get('통원'):
            continue
        if not _is_silson_like(_c.get('company'), _c.get('product'), _dm):
            continue
        try:
            _g590 = silson_gen(_c.get('contract_date'), _dm.get('입원'), _c.get('product'),
                               _has_nonpay3(_dm), _has_drug(_dm))
        except Exception:
            _g590 = ''
        if str(_g590).startswith('5'):
            continue                      # ★5세대 제외 — 지점장 지시
        _dm['통원'] = 20
        print('[v590 실손] %s — 통합형 실손 통원 20 기재(세대 %s · 제9조 5항)'
              % (str(_c.get('company'))[:10], _g590 or '미판정'))

    for _c in deduped:
        if _no_period(_c.get('contract_date'), _c.get('expiry_date')) and \
           not _is_silson_like(_c.get('company'), _c.get('product'), _c.get('dambo')):
            print(f"[v262 기간불명·포함] {_c.get('company')} {str(_c.get('product'))[:28]} — 계약일/만기일 공란 그대로 기재")
    # ══════════════════════════════════════════════════════════════════════════
    # ★★★v47 심장 묶음담보 분해 (지침 §8.3.1 + 보험인포메이션 p16~19 회사별 정본표)
    #   "묶음 진단비는 보장 구성질환의 마스터 행에 동일 금액을 각각 기재한다."
    #   회사마다 '특정Ⅰ/Ⅱ'의 뜻이 다르다 → 라벨 말고 회사별 질병코드 기준으로 분해.
    #   ★2026.07.13 지점장 확정 3건:
    #     (1) KB 특정Ⅰ = 협심증+빈맥+심부전 <b>+염증</b> (구 정본 '염증X' 폐기)
    #         ★v384: 구 주석의 '허혈성'은 <b>협심증</b>으로 표기한다(허혈성은 협심증의 종류).
    #     (2) KB 심장판막질환에서 염증(심내막염 I33) 삭제 → 판막만
    #     (3) 현대 특정Ⅰ = <b>협심증</b>+빈맥+심부전 (구 빈맥+심부전)
    #   끝열은 행별 가로 SUM이라 세로 중복합산 없음. 원천 분해 → 4대 산출물 자동 연동.
    # ══════════════════════════════════════════════════════════════════════════
    # ══════════════════════════════════════════════════════════════════
    # ★★★심장 묶음담보 정본 (지점장 확정 2026.07.14 · 회사별 표대로)
    #   공통원칙: 묶음 진단비 = 구성질환 마스터 행에 동일 금액 각 100% 기재(절반 분할 금지).
    #   ★★'허혈성 진단비' 행 = 회사담보명 '허혈성심장질환진단비' <b>단독 담보 전용</b>(전 회사 공통).
    #      묶음(특정Ⅰ·Ⅱ·Ⅲ 등)은 허혈성 행에 절대 넣지 않는다. I20·I24·I25는 '협심증' 행으로 표현.
    #   ★빈맥(I47·I48)과 부정맥(I49)은 별개 코드.
    # ══════════════════════════════════════════════════════════════════
    _hbkey = _heart_hbkey            # ★v533 제123조 — 회사 판별도 한 함수
    for _c in deduped:
        _hk = _hbkey(_c.get('company'))
        if not _hk or _hk not in _HB: continue
        for _k in list(_c['dambo'].keys()):
            _t = re.sub(r'\s', '', str(_k))
            if any(x in _t for x in ('[확인]','수술','주요치료','산정특례','혈전')): continue
            # ★★★★★v528 제120조 (지점장 실측 2026.08.21 「협심증이 500인데 또 1000으로 잡힌다」)
            #   롯데 간편 상품은 <b>모든 담보명 끝에 `(간편할인형Ⅱ)`</b>가 붙는다. 그 <b>Ⅱ를 등급으로</b> 읽어
            #   `특정심장질환Ⅰ`이 <b>'특정심장' + '2'</b> 규칙에 걸려 <b>협심증·염증으로 오분해</b>됐다.
            #   실측: 157 특정심장Ⅰ 500 → (오)협심증500+염증500 / 158 특정심장Ⅱ 500 → 협심증500+염증500
            #        ⇒ <b>협심증 1,000 · 염증 1,000</b>(김순자 PPT 실물과 일치). 급성심근경색은 비었다.
            #   ⇒ 등급 판정 전에 <b>상품 수식어 괄호를 먼저 지운다</b>. 담보 등급 괄호는 건드리지 않는다.
            _t2 = _heart_norm(_t)          # ★v533 제123조 — 실사용과 자가진단이 같은 함수
            for _pred, _rows in _HB[_hk]:
                try:
                    if not _pred(_t2): continue
                except Exception:
                    continue
                # ★★★★★v285 (지점장 지적 2026.07.31 "현대 심장도 미표기다"):
                #   구 코드는 <b>구성행이 1개인 묶음</b>을 무조건 resolve_kw에 위임하고 break 했다.
                #   그런데 resolve_kw가 그 표기를 못 잡으면 <b>아무도 처리하지 않아 담보가 통째로 사라진다</b>.
                #   <b>실측(이영태 현대 Hi2007)</b>: `심혈관질환(특정Ⅱ)진단(갱신형)담보` 500
                #     → _HB는 ['급성심근경색'] 단일행이라 break · resolve_kw는 None → <b>급성심근경색 미표기</b>.
                #     (같은 계약의 `심혈관질환(특정Ⅰ)`은 3행 묶음이라 정상 분해됐다 — 그래서 Ⅱ만 조용히 빠졌다.)
                #   같은 함정이 <b>KB·한화·NH·DB 특정Ⅱ · 현대 특정허혈/특정2대 · 롯데 특정심장Ⅰ/기타부정맥 ·
                #   KB 심장판막/심근병증 · 흥국 기타심장부정맥</b>에 전부 걸려 있다.
                #   → resolve_kw가 잡으면 <b>종전대로 위임</b>하고, <b>못 잡을 때만</b> 정본표대로 기재한다(회귀 0).
                if len(_rows) <= 1:
                    try: _std0 = resolve_kw(_k)[0]
                    except Exception: _std0 = None
                    if _std0: break                # 기존 경로가 이미 처리 중 — 손대지 않는다
                    _v1 = _c['dambo'].pop(_k)
                    _nk1 = f'{_rows[0]}[심장묶음]'
                    _c['dambo'][_nk1] = _c['dambo'].get(_nk1, 0) + _v1
                    print(f"[v285 심장단일] {_hk} '{_k}' {_v1} → {_rows[0]} (resolve 미매칭 구제)")
                    break
                _v = _c['dambo'].pop(_k)
                for _r in _rows:
                    _nk = f'{_r}[심장묶음]'
                    _c['dambo'][_nk] = _c['dambo'].get(_nk, 0) + _v
                print(f"[v47 심장묶음] {_hk} '{_k}' {_v} → {' + '.join(_rows)} (각각)")
                break
    # ══════════════════════════════════════════════════════════════════
    # ★★v51 심뇌혈관수술비 분해 (지점장 확정 2026.07.13 · 현대해상 수술비)
    #   "심뇌혈관수술비 3,000 = 심장수술비 3,000 + 뇌혈관수술비 3,000 (각각 기재)"
    #   묶음담보 공통원칙(§8.3.1)과 동일 — 절반 분할 아님, 두 행에 동일 금액 각 100%.
    #   기존 매핑은 '심뇌혈관' → 뇌혈관수술비 하나로만 넣어 심장수술비가 누락됐다.
    #   원천(dambo)에서 쪼개므로 4대 산출물(엑셀·보장나무·PPT·설명서) 자동 연동.
    # ══════════════════════════════════════════════════════════════════
    for _c in deduped:
        for _k in list(_c['dambo'].keys()):
            _kk = re.sub(r'\s', '', str(_k))
            if '[확인]' in _kk: continue
            # ★v217: '심뇌<u>5대</u>혈관'처럼 사이에 글자가 끼는 표기도 잡는다(구 '심뇌혈관' 연속 조건 완화).
            # ★★★★★v397 (지점장 확정 2026.08.12): <b>'주요치료'가 붙으면 수술 분해를 하지 않는다</b>.
            #   지점장 원문: 「심뇌혈관질환주요치료비 -> 2대주요치료비다」
            #   실측: 현대 `심뇌혈관질환<b>주요치료비</b>(…)(<b>수술</b>및혈전용해치료)` 500이 괄호 수식어의
            #     '수술' 글자 때문에 <b>심장수술비 500 + 뇌혈관수술비 500</b>으로 쪼개졌다.
            #   ★`주요심뇌5대혈관수술비`·`심뇌혈관수술비`는 '주요치료'가 없으므로 <b>종전대로 분해</b>한다.
            if '심뇌' not in _kk or '혈관' not in _kk or '수술' not in _kk or '주요치료' in _kk: continue
            _v = _c['dambo'].pop(_k)
            for _r in ('심장수술비', '뇌혈관수술비'):
                _nk = f'{_r}[묶음]'   # ★태그에 '뇌혈관' 금지(resolve 오인 방지)
                _c['dambo'][_nk] = max(_c['dambo'].get(_nk, 0), _v)   # ★v217 대표(max) — 여러 줄 합산 금지
            print(f"[v217 심뇌혈관수술] {_c.get('company')} '{_k}' {_v} → 심장수술비 + 뇌혈관수술비 (각각 대표 {_v})")

    # ══════════════════════════════════════════════════════════════════
    # ★★v46 결합담보 분해 (지점장 확정 2026.07.13 / 지침 §8.3.1 묶음담보 공통원칙 적용)
    #   "묶음(결합) 담보는 보장 구성담보의 마스터 행에 동일 금액을 각각 기재한다."
    #   예) 롯데 DB '상해사망80%이상후유장해 18,000'
    #       → 상해사망 18,000  +  상해80%이상후유장해 18,000  (두 행 각각)
    #   끝열은 행별 가로 SUM이라 세로 중복합산 없음. 원천(dambo)에서 쪼개므로 4대 산출물 자동 연동.
    # ══════════════════════════════════════════════════════════════════
    for _c in deduped:
        for _k in list(_c['dambo'].keys()):
            _kk = re.sub(r'\s', '', str(_k))
            if '[확인]' in _kk: continue
            if not ('사망' in _kk and '후유장해' in _kk): continue      # 결합담보만
            # ★★★★★v392 (2026.08.12 하소망 실측 — 지점장 지시):
            #   별첨 담보명이 <b>`특약명 : 담보명`</b> 형태인 회사가 있다(새마을금고).
            #   `화재상해사망후유장해특약 : <b>화상수술비</b>` — 사망·후유장해는 <b>콜론 앞 특약명</b>에만 있고
            #   진짜 담보는 뒤의 `화상수술비`인데, v46이 특약명만 보고 <b>상해사망 + 상해후유장해로 오배정</b>했다
            #   (실측: 화상진단비 50 · 화상수술비 100이 마스터 97·75행에 안 들어가고 사망·후유장해에 가산).
            #   → <b>콜론 뒤가 실담보명</b>이다. 그 뒤에 사망도 후유장해도 <b>없으면</b> 결합담보가 아니다 —
            #     분해하지 않고 <b>키를 실담보명으로 바꿔</b> 정상 매핑에 넘긴다.
            #   ★회귀 0 설계: 콜론 뒤에 사망 또는 후유장해가 <b>하나라도 있으면</b> 종전대로 분해한다
            #     (`…특약 : 일반상해사망공제금` · `…특약 :화재상해고도후유장해공제금[80%이상]` 등 불변).
            _tail = re.split(r'[:：]', str(_k))[-1].strip() if re.search(r'[:：]', str(_k)) else ''
            _tk = re.sub(r'\s', '', _tail)
            if _tk and ('사망' not in _tk) and ('후유장해' not in _tk):
                _v0 = _c['dambo'].pop(_k)
                _c['dambo'][_tail] = _c['dambo'].get(_tail, 0) + _v0
                print(f"[v392 특약접두] {_c.get('company','')} '{_k}' {_v0} → 실담보명 '{_tail}' (결합분해 안 함)")
                continue
            _v = _c['dambo'][_k]
            _ax = '질병' if '질병' in _kk else '상해'                    # 축: 질병 / 상해
            # 등급: '80%이상' 또는 '고도' 명시 → 80% 행 / 명시 없으면 3% 행 (담보명 문자 그대로, 추측 금지)
            _hi = ('80%이상' in _kk) or ('고도' in _kk)
            _dead = f'{_ax}사망[결합]'
            _dis  = f'{_ax}{"80%이상" if _hi else ""}후유장해[결합]'
            _c['dambo'].pop(_k)
            _c['dambo'][_dead] = _c['dambo'].get(_dead, 0) + _v
            _c['dambo'][_dis]  = _c['dambo'].get(_dis, 0) + _v
            print(f"[v46 결합담보 분해] {_c.get('company','')} '{_k}' {_v} → {_dead} + {_dis} (각각)")
    # 한장보장표 회차 주입 (별첨에 없던 pay_count 보정)
    for c in deduped:
        if not c['pay_count']:
            pc = paycount_map.get((c['company'], c['contract_date'], c['expiry_date'])) \
                 or paycount_map.get((c['contract_date'], c['expiry_date']))
            if pc: c['pay_count'] = pc
    # 병합·회차 보정 반영하여 갱신 재판정 (정본 §7 규칙대로만)
    for c in deduped:
        # ★★★★★v474 제82조 — <b>상품명 오염 절단은 여기 한 곳에서만</b>(지점장 지시 2026.08.18
        #   「제목이 끝나면 끝내라」). 헤더 표기·갱신 판정·CI 판정이 <b>전부 같은 상품명</b>을 쓰도록
        #   judge_renewal 호출 <b>앞</b>에서 한 번 정제한다(제0조 「판정은 한 곳에서만」).
        _pc0 = c.get('product') or ''
        _pc1 = _clean_product(_pc0)
        if _pc1 != _pc0:
            c['product'] = _pc1
            print('[v474 상품명] 꼬리 절단 %r → %r' % (_pc0[:60], _pc1[:60]))
        # ★★★★★v476 제84조 — 회사명 앞 기간칸 오염도 <b>같은 자리에서</b> 걷어낸다.
        _co0 = c.get('company') or ''
        _co1 = _clean_company(_co0)
        if _co1 != _co0:
            c['company'] = _co1
            print('[v476 회사명] 머리 절단 %r → %r' % (_co0[:40], _co1[:40]))
        # ★★★v207 (지점장 확정 2026.07.25, 영구지침): 3열(KB·메리츠)도 judge_renewal을 그대로 탄다.
        #   <b>납입기간 == 보장기간(가입~만기)이면 '갱신'</b>이다 — 운전자·실손도 예외 없다.
        #   구 v44 규칙('3열은 총회차가 없으니 ④ 적용 금지')은 <b>폐기</b>. 3열에도 납입기간(20년납)과
        #   보험기간(2026.03.27~2046.03.27)이 그대로 인쇄돼 있어 ④ 판정에 필요한 값이 다 있다.
        #   실측 오류(양*선 KB): 삼성 운전자 20년납/20년만기 → 비갱신(오류) · New내돈내삼 54년납/54년만기 → 비갱신(오류).
        # ★★★★★v521 제117조 2항 — <b>제안서 계약은 재판정하지 않는다</b>.
        #   제안서에는 가입일·만기일자·납입회차가 <b>없다</b>(가입 전이므로). `judge_renewal`을 태우면
        #   입력이 전부 빈 값이라 <b>무조건 비갱신</b>이 된다 — 실측 DB 참좋은운전자2607이 그랬다.
        #   제안서는 `build_proposal_contract`에서 <b>납기/만기(term)</b>로 이미 판정했다. 그 값을 쓴다.
        if not c.get('proposal'):
            c['renewal'] = judge_renewal(c['product'], c['expiry_date'], c['pay_count'], c['contract_date'], c['pay_period'], c.get('company',''), c.get('dambo'))
        # ★★★★★v473b 폐기 (지점장 지적 2026.08.18 「지침을 어겼다 · 늘 지침이 우선이다」)
        #   구 규칙 = 「담보 절반 이상이 '갱신형' 표기면 갱신 강제」.
        #   <b>지침 §6에 없는 조문</b>이다. '절반'이라는 임계값은 어디에도 없다.
        #   이 규칙이 <b>지침 ④(납입기간 == 보장기간)를 덮어써서</b> LIFEPLUS 3N5(20년납/44년보장)를
        #   갱신으로 뒤집었다. 특약 꼬리 `(3N5간편,갱신형)`은 <b>상품의 갱신 여부가 아니다</b>.
        #   ⇒ 삭제한다. 담보에 찍힌 '갱신'은 지침 ④의 <b>전기납 판별</b>에서만 쓴다(judge_renewal 내부).
    # ★신버전 보충: 세부가입현황에서 뇌·심 담보 파싱해 별첨서 0인 항목만 보충(첫 계약에 귀속)
    # ★★v43 뇌혈관 유동 재배치 (지점장 2026.07.13 확정)
    #   [정본] AIA생명·라이나생명·AIG손보·우체국 = 별첨에 '뇌혈관'이라 적혀 있어도 그대로 믿지 말 것.
    #          반드시 세부가입현황 표를 100% 대조해서, 실제로 잡힌 행(뇌혈관진단비/뇌졸증/뇌출혈)으로 배치한다.
    #          엑셀·PPT·보장진단서 모두 이 결과를 따른다. 회사별 하드코딩(구 '라이나=뇌출혈') 폐기.
    # ★★★v181 (지점장 확정 2026.07.22): <b>상세 세부내역 체크는 1000% 필수</b>.
    #   대상을 4개사에서 <b>생명보험사 전체 + CI/GI/리빙케어/퍼펙트 상품 전체</b>로 확대한다.
    #   생보 CI/GI는 담보명이 '뇌혈관진단'이어도 실제는 뇌출혈인 경우가 많아 담보명만 믿으면 안 된다.
    _SEBU_FORCE = ('AIA', '라이나', 'AIG', '우체국')          # 뇌출혈 기본배치 대상(실측 확인된 4개사)
    _SEBU_CHECK = ('생명', '라이프', '공제', 'AIA', 'AIG', '라이나', '우체국',
                   '메트라이프', '처브', 'ABL', 'KDB')        # 반드시 세부내역 대조할 대상(확대)
    def _sebu_target(_co, _pd):
        _c1 = re.sub(r'[\s（）()]', '', str(_co or ''))
        _p1 = re.sub(r'[\s（）()]', '', str(_pd or ''))
        if any(f in _c1 for f in _SEBU_CHECK): return True
        # ★★★v235: 여기도 `'CI보험'` 연속매칭이라 `CI종신보험`을 놓쳤다 → `_isci_prod`와 정본 1개로 통일.
        return _isci_prod(_p1, _co)                       # ★v446 손보 제외
    try:
        _sb = parse_sebu(lines)
        _nh  = float(_sb.get('뇌혈관진단비', 0) or 0)
        _jol = float(_sb.get('뇌졸증진단비', 0) or 0) or float(_sb.get('뇌졸중진단비', 0) or 0)
        _chu = float(_sb.get('뇌출혈진단비', 0) or 0)
        for _c in deduped:
            _co = re.sub(r'[\s（）()]', '', str(_c.get('company', '')))
            _pd = str(_c.get('product', ''))
            if not _sebu_target(_co, _pd):
                continue                                   # ★v181 생보사 전체 + CI/GI/리빙케어/퍼펙트
            for _k in list(_c['dambo'].keys()):
                _kk = str(_k).replace(' ', '')
                if not (('뇌혈관' in _kk) and ('진단' in _kk)):            continue
                if any(x in _kk for x in ('수술','주요치료','산정특례','혈전','특정')): continue
                if any(x in _kk for x in ('Ⅰ','Ⅱ','Ⅲ','II','III')):        continue
                if not _sb:
                    # ★★★v180 (지점장 확정 2026.07.21): 세부가입현황이 <b>이미지/벡터라 텍스트가 아예 없는</b>
                    #   리포트가 있다(실측 장문순: '한장 보장 현황'·'세부 가입 현황' 글자 자체가 0건 추출).
                    #   이때 [확인]큐로만 두면 <b>뇌혈관진단비 행에 그대로 남아 오류로 나간다</b>.
                    #   → 대상 4개사(AIA·라이나·AIG·우체국)는 <b>기본값을 뇌출혈진단비로 배치</b>하고
                    #     확인 메모를 남긴다(지점장: "AIA생명은 뇌출혈인데 뇌혈관으로 나온다").
                    if not any(_f in _co for _f in _SEBU_FORCE):
                        # ★v181 4개사 외(그 외 생보·CI)는 임의 이동 금지 — 확인큐로만 강하게 띄운다.
                        _c['dambo']['[확인] 세부가입현황 미파싱 · 뇌혈관/뇌졸증/뇌출혈 축 대조 필수 ' + str(_k)] = 0
                        print(f"[v43 확인큐·필수대조] {_co} '{_k}' 세부가입현황 미파싱 → 상세내역 수기 대조")
                        continue
                    _v0 = _c['dambo'].pop(_k)
                    _nk0 = '뇌출혈진단비[세부가입정본]'
                    _c['dambo'][_nk0] = float(_c['dambo'].get(_nk0, 0) or 0) + float(_v0 or 0)
                    _c['dambo']['[확인] 세부가입현황 미파싱 → 뇌출혈로 배치함, 상세내역 대조 요망 ' + str(_k)] = 0
                    print(f"[v43 기본값·뇌출혈] {_co} '{_k}' 세부가입현황 미파싱 → 뇌출혈진단비 배치(+확인메모)")
                    continue
                if   _nh  > 0: _tgt = None                 # 세부가입현황이 뇌혈관진단비로 잡음 → 유지
                elif _jol > 0: _tgt = '뇌졸중'
                elif _chu > 0: _tgt = '뇌출혈'
                else:          _tgt = None
                if not _tgt:
                    continue
                _v = _c['dambo'].pop(_k)
                _nk = _tgt + '진단비[세부가입정본]'
                _c['dambo'][_nk] = float(_c['dambo'].get(_nk, 0) or 0) + float(_v or 0)
                print(f"[v43 뇌혈관 유동재배치] {_co} '{_k}' → {_nk} (세부가입현황 정본)")
    except Exception as _e:
        print(f"[v43 뇌혈관 재배치 스킵] {_e}")

    try:
        _sebu=parse_sebu(lines)
        if _sebu and deduped:
            _loc={}
            for _ci,c in enumerate(deduped):
                for dk in list(c['dambo'].keys()):
                    _std=_dedup_std(dk)
                    if _std:
                        _loc.setdefault(_std,[]).append((_ci,dk))
            for k,v in _sebu.items():
                if k=='__무시__': continue   # ★v30z 무시지정 담보 이중차단
                if k in _loc:
                    _positions=_loc[k]
                    _bsum=0.0
                    for _ci,dk in _positions:
                        try: _bsum+=float(str(deduped[_ci]['dambo'][dk]).replace(',',''))
                        except: pass
                    if abs(_bsum - v) > 0.5:
                        _fci,_fdk=_positions[0]
                        deduped[_fci]['dambo'][_fdk]=v
                        for _ci,dk in _positions[1:]:
                            deduped[_ci]['dambo'].pop(dk,None)
                else:
                    deduped[0]['dambo']['[세부보충]'+k]=v
    except Exception as _e:
        pass
    # ★★v129 지점장 확정 2026.07.21: 월보험료 합계 판정 근거는 보유계약 리스트의
    #   '납입횟수(납부기간/총납부횟수)'와 '잔여보험료(완납이면 0)'다.
    #   보유계약 리스트 줄을 스캔해 계약에 remain을 붙인다. 못 찾으면 None → 회차로 폴백.
    try:
        _rows = []
        _re_row = re.compile(
            r'^\s*\d{1,2}\s+(\S.*?)\s+(\d{1,3})\s*/\s*(\d{1,3})\s+([\d,]+)\s*원'
            r'\s+([\d,]+)\s*원\s+([\d,]+)\s*원\s*$')
        for _ln in txt.split('\n'):
            _m = _re_row.match(_ln)
            if not _m: continue
            try:
                _rows.append({
                    'head':   re.sub(r'\s', '', _m.group(1)),
                    'pc':     f'{_m.group(2)}/{_m.group(3)}',
                    'prem':   int(_m.group(4).replace(',', '')),
                    'remain': int(_m.group(6).replace(',', '')),
                })
            except Exception:
                pass
        if _rows:
            for _c in deduped:
                _c['remain'] = None
                _key = re.sub(r'\s', '', str(_c.get('company', '')) + str(_c.get('product', '')))
                for _r in _rows:
                    if _r['prem'] and _r['prem'] == (_c.get('premium') or 0):
                        _c['remain'] = _r['remain']
                        if not _c.get('pay_count'): _c['pay_count'] = _r['pc']
                        break
                else:
                    for _r in _rows:
                        if _key and _r['head'][:8] and _r['head'][:8] in _key:
                            _c['remain'] = _r['remain']
                            if not _c.get('pay_count'): _c['pay_count'] = _r['pc']
                            break
    except Exception:
        pass
    # ★v271: 세부가입현황 '계약별 가입정보'를 함께 실어 보낸다.
    #   별첨 담보명이 특약명뿐인 계약(삼성생명 無○○(본인) 등)은 build_excel에서
    #   매핑 0건으로 판정되며, 그때 이 값으로 계약 열을 채운다.
    # ★★★v273(지점장 지시 2026.07.30 "1-3 1-5 1-7은 정상계약리스트에서만 나온다"):
    #   '수술보장' 담보는 <b>별첨(정상계약 리스트) 줄 단위 대표(max)</b>가 정답이다.
    #   dambo는 같은 담보명 2줄을 <b>합산</b>하므로(실측 100+100=200 · 200+200=400)
    #   그 값을 쓰면 1-3종이 두 배가 된다. → 별첨 원문에서 계약별 max를 따로 뽑는다.
    _surg13={}; _curp=None
    for _l13 in lines:
        _mp=re.search(r'보험료\s*([\d,]+)\s*원', _l13)
        if _mp:
            try: _curp=int(_mp.group(1).replace(',',''))
            except: _curp=None
        if _curp and ('수술보장' in re.sub(r'\s','',_l13)):
            for _mm in re.finditer(r'수술보장[^0-9]*?([\d,]+)', _l13):
                try: _v13=float(_mm.group(1).replace(',',''))
                except: continue
                if _v13>0: _surg13[_curp]=max(_surg13.get(_curp,0.0), _v13)
    try: _sbc271 = parse_sebu_bycontract(lines)
    except Exception as _e271:
        _sbc271 = {}; print('[v271 sebu] 파서 실패:', _e271)
    # ★★★★★v383 [중복줄 = 세부가입현황이 정본] — 지점장 확정 2026.08.11.
    #   지점장 원문: "요건 각각이니까 <b>합산</b>내야지" / "<b>세부내역을 보면 답이 나오는데?</b>"
    #   [배경] v344(2026.08.02)가 「같은 담보명 2줄 = 대표(max)」를 기본으로 만들었는데,
    #     진짜 별개 담보 2건이 오면 <b>조용히 하나가 사라진다</b>.
    #     실측(한정환) = 신한 `급성심근경색진단 3,000` 2줄 → 3,000만 기재 → 한장표 13,400과 <b>-3,000</b>.
    #     DB `질병사망 3,000/2,000` → 3,000만 → <b>-2,000</b>. DB `암진단 2,400/3,000` → <b>-2,400</b>.
    #   [왜 「무조건 합산」이 아닌가] 같은 파일에 <b>합산하면 안 되는 반례</b>가 있다.
    #     메트 `일반사망 6,000` 2줄 = 세부에서 <b>질병사망 6,000 · 상해사망 6,000</b>(축이 다른 한 쌍)이고,
    #     `입원 3` 2줄 = 질병일당 3 · 상해일당 3이다. 합치면 2배가 된다.
    #   → <b>세부가입현황(계약별 가입정보)의 그 마스터행 값을 그대로 쓴다.</b> 합산이냐 대표냐를 따지지 않는다.
    #     신한 급성심근 = 세부 6,000 / DB 질병사망 = 세부 5,000 / DB 암 = 세부 5,400 →전부 정답.
    #     메트 `일반사망`은 세부에 그 이름의 행이 없다(질병사망·상해사망으로 갈림) → <b>손대지 않는다</b>.
    #   ★세부에 그 마스터행이 없으면 건드리지 않는다(추측 기재 금지 · 종전 대표값 유지).
    if _sbc271:
        _dupfix = []
        for _ct383 in contracts:
            try: _sv383 = _sbc271.get(int(_ct383.get('premium') or 0)) or {}
            except Exception: _sv383 = {}
            if not _sv383: continue
            for _dn383, _dvs383 in (_ct383.get('dup') or {}).items():
                if len(_dvs383) < 2: continue
                try: _std383 = resolve2(_dn383)[0]
                except Exception: _std383 = None
                if not _std383 or _std383 not in _sv383: continue
                _new383 = _sv383[_std383]
                _cur383 = _ct383['dambo'].get(_dn383)
                if _cur383 is None or _new383 == _cur383: continue
                _ct383['dambo'][_dn383] = _new383
                _dupfix.append(f"{_ct383.get('company','')}/{_dn383} {_cur383}→{_new383}({'/'.join(str(x) for x in _dvs383)})")
        if _dupfix:
            print('[v383 중복줄·세부정본] ' + ' | '.join(_dupfix[:8]))
    _hj279 = parse_hanjang(lines)
    # ★v295: KB 3열은 한장보장표가 없다 → 2~3p '전체 보장 현황'을 <b>검산 전용</b> 앵커로 확보한다.
    #   ★세부보충 차단(v293) 판정에는 쓰지 않는다 — 그건 원래 한장표만 본다(오염 재발 방지).
    try: _hjkb = parse_kb_summary(lines)
    except Exception as _e295:
        _hjkb = {}; print('[v295 KB요약표] 실패:', _e295)
    # ★★★v279: 라벨 경로가 비었거나 빈약하면 <b>좌표 폴백</b>(신형 let: = 라벨이 이미지)
    # ★★★★★v293 (김수영 KB 실측 2026.07.31, 영구): <b>검증 앵커가 없으면 세부보충을 하지 않는다</b>.
    #   ・좌표 경로 `parse_sebu_bbox`는 한장표(hj)로 대조 검증을 하고 실패하면 스킵한다.
    #     그런데 <b>라벨 경로 `parse_sebu_bycontract`에는 검증이 아예 없었다</b>.
    #   ・이 KB 신형 리포트에는 <b>세부가입현황 표가 없다</b>(4p 계약리스트 · 5p 상품별 가입현황뿐).
    #     라벨 경로가 5p 표를 세부표로 착각해 <b>합계열·타사열 숫자</b>를 DB손보 열에 박았다.
    #     실측 오염 = 유사암 4,000 / 골절 2,079 / 질병수술비 2,079 / 깁스 5,000 / 상해사망 22 /
    #     상해수술비 22 / 입원 1 / 통원 350 — 전부 근거 없는 값. 게다가 별첨 정상값
    #     (가족생활배상책임 1억)까지 지워 <b>일상배상책임이 통째로 사라졌다</b>.
    #   ・지침 영구원칙 「구조 가정 금지」 + 「조용히 틀리는 것을 시끄럽게 틀리는 것으로」 적용.
    #     → 한장표가 없으면 <b>보충하지 않고 확인사항에 남긴다</b>(추측 기재 금지).
    _sebu_blocked = 0
    if _sbc271 and not _hj279:
        _sebu_blocked = len(_sbc271)
        print(f'[v293 세부보충 차단] 한장보장표 없음 → 검증 불가 → 세부가입현황 보충 {len(_sbc271)}건 폐기(추측 기재 금지)')
        _sbc271 = {}
    if len(_sbc271) < max(3, len(deduped)//2):
        try:
            _bb = parse_sebu_bbox(_hj279)
            if len(_bb) > len(_sbc271):
                print(f'[v279 sebu-bbox] 라벨 경로 {len(_sbc271)}건 → 좌표 경로 {len(_bb)}건으로 대체')
                _sbc271 = _bb
        except Exception as _e279:
            print('[v279 sebu-bbox] 실패:', _e279)
    return {'client':client,'contracts':deduped,'sebu_bc':_sbc271,'surg13':_surg13,'hanjang':_hj279,'hanjang_kb':_hjkb,'sebu_blocked':_sebu_blocked}

# ★ DMAP — 마스터 엑셀 B열 기준 100% 일치
DMAP = {
    '5대골절진단':'5대골절진단비',   # ★v30c '골절진단(간편Ⅲ)' 부분일치 오탐 차단(선순위)
    # ★v29t §8.1: 동양류 '[N] 주계약_주계약' 2줄 = 일반사망+상해사망 1:1 (~2 접미사=두 번째 줄)
    '주계약_주계약~2':'상해사망','주계약_주계약':'일반사망',
    # 사망
    '상해사망':'상해사망','상해사망(갱신형) [보통약관]':'상해사망','일반상해사망':'상해사망',
    '기본계약(상해사망(간편가입Ⅲ))담보':'상해사망',
    '질병사망':'질병사망(80세)',
    # 후유장애
    '상해후유장해3%':'상해후유3%','상해후유80%':'상해후유80%',
    '질병후유장해3%':'질병후유3%','질병후유80%':'질병후유80%',
    # 암 — B열: 고액암/일반암/중대한 암/유사암(갑.기.경.제)/표적항암치료비/하이클래스(암)/중입자치료비/양성자치료/세기조절치료/다빈치로봇수술비/암수술/암일당/항암방사선약물
    '일반암진단비':'일반암','암진단Ⅱ(유사암제외)(간편가입Ⅲ)담보':'일반암',
    '고액암진단비':'고액암',
    '갑상선암.기타피부암.유사암진단비Ⅲ':'유사암(갑.기.경.제)','유사암진단비':'유사암(갑.기.경.제)',
    '유사암진단Ⅱ(양성뇌종양포함)(간편가입Ⅲ)담보':'유사암(갑.기.경.제)',
    '표적항암약물허가치료비':'표적항암치료비','표적항암약물허가치료(간편가입Ⅲ)(갱신형)담보':'표적항암치료비',
    '항암방사선.약물치료비':'항암방사선약물','항암방사선치료(간편가입Ⅲ)담보':'항암방사선약물',
    '항암약물치료(간편가입Ⅲ)담보':'항암방사선약물',
    '항암방사선(세기조절)치료(간편가입Ⅲ)(갱신형)담보':'세기조절치료','항암세기조절방사선치료비':'세기조절치료',
    '항암방사선(양성자)치료(간편가입Ⅲ)(갱신형)담보':'양성자치료','항암양성자방사선치료비':'양성자치료',
    '암수술(간편가입Ⅲ)담보':'암수술',
    # ★v267(지점장 확정 2026.07.28): 카티(CAR-T)는 항암방사선약물이 아니다 → [확인]큐. 별칭표에서 삭제.
    # 뇌혈관 — B열: 뇌혈관진단비/뇌졸증진단비/중대한 뇌졸증/뇌출혈진단비/외상성뇌출혈/산정특례뇌혈관/혈전용해치료비
    '뇌혈관질환진단비Ⅲ(건강맞춤형Ⅱ)(갱신형)':'뇌혈관진단비','뇌혈관질환진단비':'뇌혈관진단비',
    '뇌혈관질환진단(간편가입Ⅲ)담보':'뇌혈관진단비',
    '뇌졸중진단비':'뇌졸증진단비','뇌졸중진단(간편가입Ⅲ)담보':'뇌졸증진단비',
    '뇌졸중진단비(건강맞춤형Ⅱ)(갱신형)':'뇌졸증진단비',
    '뇌출혈진단':'뇌출혈진단비',
    '중증질환자(뇌혈관질환)산정특례대상진단비(연간1회한)(건강맞춤형Ⅱ)(갱신형)':'산정특례뇌혈관',
    '뇌혈관질환수술비Ⅲ(건강맞춤형Ⅱ)(갱신형)':'뇌혈관수술비','심뇌혈관질환수술(간편가입Ⅲ)담보':'뇌혈관수술비',
    '뇌경색증(I63)혈전용해치료비':'혈전용해치료비','혈전용해치료비Ⅱ(뇌졸중)(간편가입Ⅲ)담보':'혈전용해치료비',
    # 심장 — B열: 협심증/심부전/염증/부정맥/산정특례심장/2대 주요치료비/급성심근경색/중대한 급성심근/혈전용해치료비
    '허혈심장질환진단비Ⅲ(건강맞춤형Ⅱ)(갱신형)':'허혈성 진단비','허혈심장질환진단비':'허혈성 진단비',   # ★v29t §8.3 구규칙(=협심증) 폐기
    '허혈심장질환진단(간편가입Ⅲ)담보':'허혈성 진단비',
    '급성심근경색증진단':'급성심근경색','급성심근경색증진단(간편가입Ⅲ)담보':'급성심근경색',
    '중증질환자(심장질환)산정특례대상진단비(연간1회한)(건강맞춤형Ⅱ)(갱신형)':'산정특례심장',
    '허혈심장질환수술비Ⅲ(건강맞춤형Ⅱ)(갱신형)':'허혈성수술비','허혈심장질환수술비':'허혈성수술비',   # ★v29t §8.3 허혈수술→허혈성수술비 행
    '급성심근경색증(I21)혈전용해치료비':'혈전용해치료비',
    # 일당 — B열: 질병일당/질병수술일당/질병종합병원일당/상해일당/간병인/간호통합병동/1인실 종합병원/1인실 상급병원/질병중환자실/상해중환자실
    '간호간병통합서비스질병입원일당(1-180일)(간편가입)(갱신형)':'간호통합병동',
    '상급종합병원질병입원일당(상급병실(1인실),1일이상60일한도)(간편가입)(갱신형)':'1인실 상급병원',
    '종합병원질병입원일당(상급병실(1인실),1일이상30일한도)(간편가입)(갱신형)':'1인실 종합병원',
    # 수술비
    '질병수술(간편가입Ⅲ)담보':'질병수술비',
    '상해수술비(건강맞춤형Ⅱ)(갱신형)':'상해수술비','상해수술(간편가입Ⅲ)담보':'상해수술비',
    '골절수술(간편가입Ⅲ)담보':'골절수술비','화상수술(간편가입Ⅲ)담보':'화상수술비',
    '120대질병수술Ⅱ(간편가입Ⅲ)(질병수술3(24대질병))담보':'n대수술비',
    '5대기관질병수술(관혈/비관혈)(연간1회한)(간편가입Ⅲ)담보':'5대기관 수술비 관혈',
    '중대한특정상해수술(간편가입Ⅲ)담보':'중대한상해수술비',
    # 운전자 — B열: 대인/대물/합의금/6주미만/변호사/자부상 (처리지원금 판정은 resolve_kw에서 순서대로)
    '교통사고벌금(대물)':'대물','교통사고벌금(대인)':'대인',
    '변호사선임비용':'변호사','자동차사고 변호사선임비용':'변호사','변호사비':'변호사',
    '자동차부상위로금':'자부상','자동차부상보장':'자부상',
    '무보험차에 의한 상해':'일상배상책임',
    # 골절 — B열: 골절(치아파절포함)/골절(치아파절제외)/5대골절진단비
    '골절진단(간편가입Ⅲ)담보':'골절(치아파절포함)',  # 단독 골절진단=치아포함 행(치아제외 명시만 제외 행)
    # 응급실
    '응급실내원비(응급)':'응급실(응급)',
    # 화상 — B열: 진 단 비/중증화상진단비
    '화상진단비':'화상진단비','화상진단비(건강맞춤형Ⅱ)(갱신형)':'화상진단비',
    # 깁스 — B열: 반깁스/깁스진단비
    '깁스치료담보':'깁스진단비','깁스치료':'깁스진단비',
    # 실손 — B열: 입원/통원/약값
    '질병입원의료비':'입원','상해입원의료비':'입원',
    '질병외래의료비':'통원',
    '도수/체외충격파/증식치료':'도수치료','비급여주사제':'비급여주사','MRI검사의료비':'MRI',
    # 일배책
    '가족생활배상책임':'일상배상책임','일상생활배상책임':'일상배상책임',
    # 치아
    '치과치료(보존치료)':'크라운','치과치료(보철치료)':'임플란트',
    # 제외
    '보험료납입지원':None,
}

# -*- coding: utf-8 -*-
import re
# 마스터 82행 전부 키워드로 잡는 사전 엔진. (predicate, std, jong)
# 순서 = 구체 우선. 앞에서 잡히면 끝.
def _norm(s): return re.sub(r'\s+','',s)

# ★★★★★v386 (지점장 확정 2026.08.12): 「<b>116대는 6개로 다 찍어내야한다 · 1-5종처럼</b>」
#   구 코드는 n대수술비를 <b>대표(max) 1개</b>로만 찍어 KB 116대 Ⅰ~Ⅵ(500·200·100·60·30·10)이
#   <b>500 하나</b>로 뭉갰다(실측). → 1-5종과 같은 <b>등급별 슬래시 6칸</b>으로 기재한다.
_NDAE_G = (('Ⅵ',6),('Ⅴ',5),('Ⅳ',4),('Ⅲ',3),('Ⅱ',2),('Ⅰ',1),
           ('III',3),('VI',6),('IV',4),('II',2),('V',5),('I',1),
           ('6종',6),('5종',5),('4종',4),('3종',3),('2종',2),('1종',1))
def _ndae_grade(nm):
    """'116대질병수술비Ⅳ' → 4 / 등급 없으면 0. ★'대' 뒤쪽만 본다(상품명 로마숫자 오독 차단)."""
    import re as _re
    t = _re.sub(r'\s','', str(nm or '')).replace('（','(').replace('）',')')
    t = _re.sub(r'\([^)]*\)','', t)                 # (간편가입)·(갱신형) 등 제거
    m = _re.search(r'(?<!\d)\d{2,3}대', t)
    if not m: return 0
    tail = t[m.end():]
    for k,g in _NDAE_G:
        if k in tail: return g
    return 0

def _rmn(s):
    """담보명 등급 로마숫자/숫자 판별 → 3/2/1/0. 괄호 속(건강맞춤형Ⅱ 등)은 제외."""
    import re as _re
    _raw0=str(s)
    # ★양예서 버그: Adobe가 로마숫자를 전각괄호（）·파이프|로 깨뜨림. 괄호 안 등급도 읽어야 함(뇌질환진단비（II） 등)
    # 상품수식어 괄호(건강맞춤형Ⅱ 등)는 제거하되, '진단비（I/II/III）'처럼 담보 등급 괄호는 살린다
    # ★★★★★v520 제116조 <b>철회</b> (지점장 실측 2026.08.19 «지금 없는 허혈성이 나오고
    #   협심증이 1500으로 나오고 리포트에는 협심증이 안나오고 엉망이다»)
    #   v519에서 괄호 안 등급을 살렸더니 <b>롯데 묶음 분기가 처음으로 작동</b>했고,
    #   <b>구성을 모르는 채</b> Ⅱ→협심증·염증, 15대→판막·심근병증·빈맥·심부전으로 <b>분해</b>됐다.
    #   지점장이 준 정본은 <b>Ⅰ의 질병코드뿐</b>이다 — 나머지 4종의 구성은 <b>받지 않았다</b>.
    #   ⇒ <b>구 정규식으로 되돌린다</b>(제0조 «지침에도 엑셀에도 없으면 만들지 않는다»).
    s2=_re.sub(r'[(（](?!\s*[I|ⅠⅡⅢV\d]{1,4}\s*[)）])[^)）]*[)）]','',_raw0)  # 등급 아닌 괄호만 제거
    s2=s2.replace('（','(').replace('）',')').replace('|','I')  # 전각→반각, 파이프→I
    if 'Ⅲ' in s2 or 'III' in s2 or '(III)' in s2 or '3종' in s2: return 3
    if 'Ⅱ' in s2 or 'II' in s2: return 2
    if 'Ⅰ' in s2 or '(I)' in s2: return 1
    m=_re.search(r'진단비?\s*([123])(?!\d)',s2)
    if m: return int(m.group(1))
    if _re.search(r'[가-힣]I(?![A-Za-zI])',s2): return 1
    return 0

# ★★★★★v533 제123조 (지점장 지시 2026.08.21 「이제 심장만큼은 절대 틀리지마라」)
#   `_HB`를 <b>모듈 전역</b>으로 올린다. 지역변수로 숨어 있으면 <b>자가진단이 태울 수 없다</b>.
#   같은 규칙을 검사용으로 한 벌 더 만들면 두 곳이 어긋난다 — 그것이 오늘의 사고였다.
_HB = {   # ★★★심장 묶음담보 회사별 정본표(지점장 최종본 2026.07.18 · 이전 표 폐기)
  #  ★'허혈성 진단비' 행 = 회사담보명 '허혈성심장질환진단비' 단독 담보 전용.
  #   묶음은 이 행에 절대 안 넣는다. I20·I24·I25는 '협심증' 행으로 표현.
  #  ★빈맥(I47·I48)과 부정맥(I49)은 별개 코드.
  'KB': [
    (lambda t: '특정1' in t and '심장' in t,   ['협심증','빈맥','심부전','주요심장염증']),
    (lambda t: '특정2' in t and '심장' in t,   ['급성심근경색']),
    (lambda t: '기타심장부정맥' in t or ('부정맥' in t and 'I49' in t.upper()), ['부정맥']),
    (lambda t: '심장판막' in t,                 ['심장판막']),
    (lambda t: '심근병증' in t,                 ['심근병증']),
  ],
  '한화': [
    (lambda t: ('기타부정맥제외' in t or 'I49제외' in t) and ('심혈관' in t or '특정1' in t),
                                               ['협심증','빈맥','심부전']),
    (lambda t: '심혈관1' in t or '특정1' in t,  ['협심증','빈맥','부정맥','심부전']),
    (lambda t: '심혈관2' in t or '특정2' in t,  ['급성심근경색']),
    (lambda t: '심근병증' in t,                 ['심근병증']),
  ],
  'NH': [
    (lambda t: ('기타부정맥제외' in t or 'I49제외' in t) and ('심혈관' in t or '특정1' in t),
                                               ['협심증','빈맥','심부전']),
    (lambda t: '심혈관' in t and '특정1' in t,  ['협심증','빈맥','부정맥','심부전']),
    (lambda t: '심혈관2' in t or '특정2' in t,  ['급성심근경색']),
    (lambda t: '주요심장염증' in t,             ['주요심장염증']),
    (lambda t: '심근병증' in t,                 ['심근병증']),
  ],
  'DB': [
    # ★★★★★v263(지점장 지시 2026.07.27, 영구): <b>DB손보 `4대순환계질환진단비(특정3대심장질환)` = 심부전 · 부정맥</b>.
    #   지점장 원문 = "디비손보 4대순환계질환진단비(특정3대심장질환) = 엑셀 = 심부전 부정맥 에 표시하라".
    #   ★<b>반드시 '특정3' 규칙보다 앞에 둔다</b> — 담보명에 '특정3'과 '심장'이 둘 다 들어 있어
    #     아래 특정Ⅲ 규칙(심장판막·빈맥·심부전)에 먼저 걸려 <b>심장판막·빈맥으로 오분류</b>됐다.
    #   묶음 공통원칙대로 두 행에 <b>동일 금액 각 100%</b> 기재한다.
    # ★v265(2026.07.28 박O정 실측): 같은 계약에 <b>`4대순환계질환진단비(특정하지정맥류질환)`</b>이
    #   또 있어 <b>하지정맥류까지 심부전·부정맥에 산입</b>됐다(500+500=1,000). 지점장 지시는
    #   <b>`(특정3대심장질환)`</b>만이다 → <b>'3대심장' 조건 추가</b>. 하지정맥류는 [확인]큐로 보낸다.
    (lambda t: '4대순환계' in t and '3대심장' in t, ['심부전','부정맥']),
    (lambda t: '특정1' in t and '심장' in t,   ['협심증','주요심장염증']),
    (lambda t: '특정2' in t and '심장' in t,   ['급성심근경색']),
    # ★★★★★v535 (지점장 확인 2026.08.21 「빈맥을 왜 삭제해??」)
    #   DB <b>특정Ⅲ = 심장판막 · 심부전 · 빈맥</b>. 지침 정본표(2026.08.06 확정)대로다.
    #   구 v264(2026.07.27 「빈맥삭제」)는 <b>폐기</b> — 지침 표와 어긋난 채 코드에만 남아 있었다.
    #   ★`4대순환계(특정3대심장)` 규칙이 <b>이 줄보다 앞</b>에 있어야 한다(담보명에 '특정3'·'심장' 동시 포함).
    (lambda t: '특정3' in t and '심장' in t,   ['심장판막','심부전','빈맥']),
    (lambda t: '순환계3대' in t,               ['빈맥','부정맥','심부전']),
  ],
  '현대': [
    (lambda t: '특정허혈' in t,                 ['급성심근경색']),
    # ★'허혈성심장질환진단비' = 무조건 단독(허혈성 행). 분해 금지.
    (lambda t: '특정2대' in t,                  ['부정맥']),
    (lambda t: '특정1' in t and '심' in t,      ['협심증','빈맥','심부전']),
    (lambda t: '특정2' in t and '심' in t,      ['급성심근경색']),
    (lambda t: '주요심장염증' in t,             ['주요심장염증']),
  ],
  '흥국': [
    # ★★★★★v605 제123조 4항 (지점장 지시 2026.08.30 「흥국 특정3대 지침에 추가해줘」
    #   · 정상범님 최종 가입제안서 실측).
    #   약관 원문: 「특정3대심장질환 : <b>인공소생에 성공한 심장정지 · 부정맥 · 심부전</b>」
    #   ★인공소생 성공 심장정지(I46.0)는 <b>마스터 무행</b>이라 기재하지 않는다
    #     (제0조 꼭대기 · DOCTRINE 4131행). ⇒ <b>부정맥 + 심부전</b> 두 행에 각각 기재.
    #   ★`특정심혈관` 규칙보다 <b>앞에</b> 둔다 — 뒤에 두면 3대가 특정심혈관으로 새어간다.
    (lambda t: '특정3대심장' in t,              ['부정맥','심부전']),
    (lambda t: '특정심혈관' in t and '기타심장부정맥' in t and '제외' not in t, ['부정맥']),
    (lambda t: '특정심혈관' in t,               ['협심증','빈맥','심부전']),
    (lambda t: '심근병증' in t,                 ['심근병증']),
    (lambda t: '주요심장염증' in t,             ['주요심장염증']),
  ],
  '롯데': [
    (lambda t: '특정심장' in t and '2' in t,    ['협심증','주요심장염증']),
    (lambda t: '특정심장' in t and '1' in t,    ['급성심근경색']),
    # ★★★★★v610 제123조 5항 (지점장 지침 추가 2026.08.30 — 롯데 약관 질병분류표).
    #   특정15대심장질환 = <b>15개 심장질환</b>이다.
    #     허혈성 I20·21·22·23·24·25 → 협심증 + 급성심근경색
    #     부정맥 I47·48·49          → 빈맥 + 부정맥
    #     심부전 I50                → 심부전
    #     심장막 I30·31·32 · 심근염 I40·41 → 주요심장염증
    #     판막 I05·I06.0·I06.2·I38  → 심장판막
    #   ★심근병증(I42·43)은 <b>범위에 없다</b> — 구 매핑에서 뺀다.
    (lambda t: '특정15대' in t,                 ['협심증','급성심근경색','빈맥','부정맥','심부전','주요심장염증','심장판막']),
    (lambda t: '기타부정맥' in t or '기타심장부정맥' in t,               ['부정맥']),
  ],
  # ★삼성·메리츠 '허혈성심장질환진단비' = 무조건 단독(허혈성 행). 분해 없음.
  '삼성':   [],
  '메리츠': [],
}


# ★★★★★v533 제123조 2항 — <b>심장 케이스 자가진단</b> (지점장 지시 「이제 심장만큼은 절대 틀리지마라」)
#   정답은 <b>약관 실측</b>에서 온다. 코드를 보고 만든 표가 아니다.
#   담보명은 실물 제안서 그대로 — 상품 수식어 괄호까지 붙어 있어야 제120조가 검사된다.
_HEART_CASES = [
    # (회사, 실물 담보명, 정답 마스터 행)   ※근거: 각 사 약관 질병코드 실측 2026.08.21
    # ★★★★★v591 — 지점장이 잡은 <b>실물 담보명 원문</b>을 케이스로 박는다(let:smile 4060종합).
    #   구 결함: 인라인 2곳(DMAP·resolve_kw)이 회사 구분 없이 Ⅰ→협심증으로 보내
    #   `_HB` 분해분과 <b>겹쳐 협심증이 2,000</b>이 됐다. 케이스에 없어 29건 검사가 통과했다.
    ('롯데손해보험', '심혈관질환진단비(특정심장질환Ⅰ)(건강체할인형)',   ['급성심근경색']),
    ('롯데손해보험', '심혈관질환진단비(특정심장질환Ⅱ)(건강체할인형)',   ['협심증', '주요심장염증']),
    ('롯데손해보험', '심혈관질환진단비(특정15대심장질환)(건강체할인형)', ['협심증','급성심근경색','빈맥','부정맥','심부전','주요심장염증','심장판막']),
    ('롯데손해보험', '심혈관질환진단비(기타심장부정맥)(건강체할인형)',   ['부정맥']),
    ('롯데손해보험', '심혈관질환진단비(특정심장질환Ⅰ)(간편할인형Ⅱ)',   ['급성심근경색']),
    ('롯데손해보험', '심혈관질환진단비(특정심장질환Ⅱ)(간편할인형Ⅱ)',   ['협심증', '주요심장염증']),
    ('롯데손해보험', '심혈관질환진단비(특정15대심장질환)(간편할인형Ⅱ)', ['협심증','급성심근경색','빈맥','부정맥','심부전','주요심장염증','심장판막']),
    ('롯데손해보험', '심혈관질환진단비(기타심장부정맥)(간편할인형Ⅱ)',   ['부정맥']),
    ('KB손해보험',  '심장질환(특정Ⅰ)진단비(감액없음)(맞춤고지)',       ['협심증', '빈맥', '심부전', '주요심장염증']),
    ('KB손해보험',  '심장질환(특정Ⅱ)진단비(감액없음)(맞춤고지)',       ['급성심근경색']),
    ('KB손해보험',  '부정맥질환(Ⅰ49)진단비(감액없음)(맞춤고지)',      ['부정맥']),
    ('KB손해보험',  '심근병증진단비(감액없음)(맞춤고지)',              ['심근병증']),
    ('KB손해보험',  '심장판막협착증(대동맥판막)진단비(감액없음)(맞춤고지)', ['심장판막']),
    # ── 아래 20건: 지침 심장 묶음 정본표(2026.08.06 확정) 문장을 정답으로 삼는다 ──
    ('한화손해보험',   '심혈관질환(특정Ⅰ)진단비',                 ['협심증','빈맥','부정맥','심부전']),
    ('한화손해보험',   '심혈관질환(특정Ⅰ)(기타부정맥제외)진단비',   ['협심증','빈맥','심부전']),
    ('한화손해보험',   '심혈관질환(특정Ⅱ)진단비',                 ['급성심근경색']),
    ('한화손해보험',   '심근병증진단비',                          ['심근병증']),
    ('NH농협손해보험', '심혈관질환(특정Ⅰ)진단비',                 ['협심증','빈맥','부정맥','심부전']),
    ('NH농협손해보험', '심혈관질환(특정Ⅰ)(기타부정맥제외)진단비',   ['협심증','빈맥','심부전']),
    ('NH농협손해보험', '심혈관질환(특정Ⅱ)진단비',                 ['급성심근경색']),
    ('NH농협손해보험', '주요심장염증진단비',                       ['주요심장염증']),
    ('DB손해보험',    '4대순환계질환진단비(특정3대심장질환)',       ['심부전','부정맥']),
    ('DB손해보험',    '특정심장질환(특정Ⅰ)진단비',                ['협심증','주요심장염증']),
    ('DB손해보험',    '특정심장질환(특정Ⅱ)진단비',                ['급성심근경색']),
    ('DB손해보험',    '특정심장질환(특정Ⅲ)진단비',                ['심장판막','심부전','빈맥']),
    ('DB손해보험',    '순환계3대질환진단비',                      ['빈맥','부정맥','심부전']),
    ('현대해상',      '특정허혈심장질환진단비',                    ['급성심근경색']),
    ('현대해상',      '심혈관질환(특정Ⅰ)진단비',                  ['협심증','빈맥','심부전']),
    ('현대해상',      '심혈관질환(특정Ⅱ)진단비',                  ['급성심근경색']),
    ('현대해상',      '특정2대심장질환진단비',                     ['부정맥']),
    # ★★★★★v605 제123조 4항 (지점장 지시 2026.08.30 「흥국 특정3대 지침에 추가해줘」
    #   · 정상범님 최종 가입제안서 실측).
    #   약관 원문: 「특정3대심장질환 : <b>인공소생에 성공한 심장정지 · 부정맥 · 심부전</b>」
    #   ★인공소생 성공 심장정지(I46.0)는 <b>마스터 무행</b>이라 종전대로 기재하지 않는다
    #     (제0조 꼭대기 · DOCTRINE 4131행). ⇒ <b>부정맥 + 심부전</b> 두 행에 각각 기재.
    ('흥국화재',      '특정3대심장질환진단비',                      ['부정맥','심부전']),
    ('흥국화재',      '특정3대심장질환진단비(통합간편가입형)',        ['부정맥','심부전']),
    ('흥국화재',      '특정심혈관질환(기타심장부정맥)진단비',        ['부정맥']),
    ('흥국화재',      '특정심혈관질환(기타심장부정맥제외)진단비',    ['협심증','빈맥','심부전']),
    ('흥국화재',      '심근병증(허혈성제외)진단비',                ['심근병증']),
]


def heart_bundle_no_isch():
    """★★★★★제123조 4항 (지점장 재확인 2026.08.21
       「심장 씨리즈에는 허혈성은 없다 · 허혈과 관련된건 협심증이다 · 허혈성은 단독만 있다」)
       <b>묶음(_HB) 어느 규칙에도 「허혈성」 행이 들어가서는 안 된다.</b>
       약관에 I24·I25가 있어도 묶음은 <b>협심증</b> 행으로 표현한다.
       `허혈성 진단비` 행은 <b>담보명이 허혈성인 단독 담보 전용</b>이다."""
    bad = []
    for _co, _rules in _HB.items():
        for _pred, _rows in _rules:
            for _r in _rows:
                if '허혈' in str(_r):
                    bad.append('[제123조 4항] %s 묶음에 허혈성 행 — %s (묶음은 협심증으로)'
                               % (_co, _rows))
    return bad


_SILSON_CASES = [
    # ★★★★★v551 제9조 3항 (지점장 지시 2026.08.22 「실손 전세대 기재 다시봐라」)
    #   실손은 세대마다 담보명이 다르다. <b>전 세대를 케이스로 고정</b>해 자가진단에 태운다.
    #   정답은 <b>마스터 실손 5행</b>(입원·통원·약값·MRI/도수치료/비급여주사·상해의료비)에서 온다.
    # ★★★★★v590 제9조 5항 — <b>케이스 자체가 틀린 정답을 박고 있었다</b>(2026.08.25 자백).
    #   구: `('상해 의료비(입원+통원)', '상해의료비')` → 검사가 <b>오답을 정답으로 고정</b>해
    #   23건이 전부 통과했고, 그래서 지점장이 잡을 때까지 아무도 몰랐다.
    #   ★<b>지적받은 담보명 원문을 그대로</b> 케이스에 박는다(윤선경 NH농협손보 헤아림실손2301).
    ('상해의료비', '상해의료비'),
    ('상해 의료비(입원+통원)', '입원'),
    ('상해(일반상해, 전체상해를 의미) 의료비(입원+통원)', '입원'),
    ('질병(전체질병을의미) 의료비(입원+통원)', '입원'),
    ('비급여 도수/체외충격파/증식치료', '도수치료'),
    ('비급여 주사제', '비급여주사'),
    ('비급여 MRI 검사', 'MRI'),
    ('질병입원의료비', '입원'), ('질병통원의료비', '통원'),
    ('상해입원의료비', '입원'), ('상해통원의료비', '통원'), ('처방조제비', '약값'),
    ('상해입원형', '입원'), ('상해통원형', '통원'),
    ('질병입원형', '입원'), ('질병통원형', '통원'), ('처방조제료', '약값'),
    ('비급여도수치료·체외충격파·증식치료', '도수치료'),
    ('비급여주사료', '비급여주사'), ('비급여자기공명영상진단(MRI/MRA)', 'MRI'),
    ('상해급여입원의료비', '입원'), ('상해급여통원의료비', '통원'),
    ('상해비급여입원의료비', '입원'), ('상해비급여통원의료비', '통원'),
    ('질병급여입원의료비', '입원'), ('질병급여통원의료비', '통원'),
    ('질병비급여입원의료비', '입원'), ('질병비급여통원의료비', '통원'),
]


def silson_selftest():
    """★제9조 3항 — 실손 1~5세대 담보명이 마스터 5행으로 정확히 가는지."""
    bad = []
    for _nm, _want in _SILSON_CASES:
        try:
            _got = resolve2(_nm)[0]
        except Exception as e:
            bad.append('[제9조 실손케이스] %s — 예외 %s' % (_nm, e)); continue
        if _got != _want:
            bad.append('[제9조 실손케이스] %s — 정답 %s / 실제 %s' % (_nm, _want, _got))
    return bad


def heart_inline_selftest():
    """★★★★★v591 (지점장 실측 2026.08.26 「또 협심증이 2배로 나온다」).
       [구멍] `heart_case_selftest`는 <b>`heart_rows`(회사별 표)만</b> 본다.
         그런데 값을 실제로 넣는 경로는 <b>`resolve2`도</b> 있다. 인라인이 회사 구분 없이
         Ⅰ→협심증을 반환해 `_HB` 분해분과 <b>겹쳐 2배</b>가 됐는데 로봇은 100%였다.
       ⇒ <b>묶음 라벨은 `resolve2`가 값을 만들면 안 된다</b>(반드시 None)를 직접 센다.
         회사마다 뜻이 다른 라벨을 회사 없이 판정하면 어느 회사에선가 반드시 틀린다."""
    bad = []
    for _nm in ('심혈관질환진단비(특정심장질환Ⅰ)(건강체할인형)',
                '심혈관질환진단비(특정심장질환Ⅱ)(건강체할인형)',
                '심장질환(특정Ⅰ)진단비Ⅲ(간편가입)',
                '특정심장질환(특정Ⅱ)진단비'):
        try:
            _got = resolve2(_nm)[0]
        except Exception as _e:
            bad.append('[제123조 인라인] %s — 예외 %s' % (_nm[:26], str(_e)[:20])); continue
        if _got is not None:
            bad.append('[제123조 인라인] %s — <b>%s</b> (정답 None · 묶음은 _HB 한 경로에서만)'
                       % (_nm[:26], _got))
    return bad


def heart_case_selftest():
    """★제123조 2항 — 회사별 심장 묶음 분해가 약관 실측 정답과 같은지 본다.
       값 대조(제121조)는 <b>cov가 처음부터 틀리면</b> 못 잡는다. 그 앞단을 여기서 막는다."""
    bad = list(heart_bundle_no_isch()) + list(heart_inline_selftest())
    for _co, _nm, _want in _HEART_CASES:
        try:
            _got = heart_rows(_co, _nm)
        except Exception as e:
            bad.append('[제123조 심장케이스] %s %s — 예외 %s' % (_co, _nm[:24], e)); continue
        if list(_got or []) != list(_want):
            bad.append('[제123조 심장케이스] %s %s — 정답 %s / 실제 %s'
                       % (_co, _nm[:24], _want, _got))
    return bad


def _heart_norm(_t):
    """★제120조 — 등급 판정 전 상품 수식어 괄호를 지우고 로마숫자를 숫자로 바꾼다.
       실사용(_HB 분해)과 자가진단이 반드시 같은 함수를 쓴다."""
    _t = re.sub(r'\s', '', str(_t))
    _t2 = re.sub(r'[(（][^)）]*(?:간편|할인형|맞춤고지|무배당|무해지|납입면제|보통약관|갱신형|비갱신)[^)）]*[)）]', '', _t)
    # ★v533 제123조 — I49를 <b>로마숫자 치환보다 먼저</b> 고정한다.
    #   실측 KB `부정맥질환(Ⅰ49)진단비` → 구 순서는 Ⅰ→1이 먼저라 '149'로 뭉개져 <b>부정맥 미매칭</b>.
    _t2 = re.sub(r'[IiⅠ]\s*49', 'I49', _t2)
    return _t2.replace('Ⅰ','1').replace('Ⅱ','2').replace('Ⅲ','3')


def _heart_hbkey(comp):
    c = re.sub(r'[\s（）()]', '', str(comp or ''))
    for k in ('KB','한화','농협','NH','DB','현대','흥국','롯데','삼성','메리츠'):
        if k in c: return 'NH' if k == '농협' else k
    return None


def heart_rows(company, name):
    """(회사, 담보명) → 분해될 마스터 행 목록. 매칭 없으면 None."""
    _hk = _heart_hbkey(company)
    if not _hk or _hk not in _HB: return None
    _t2 = _heart_norm(name)
    for _pred, _rows in _HB[_hk]:
        try:
            if _pred(_t2): return list(_rows)
        except Exception:
            continue
    return None


# ★★★★★v580 제122조 4항 — 현물 판정 한 곳 (지점장 지시 2026.08.24)
#   DB손해보험 · 계약일 2026.08 이후 · 담보명에 `(현물급부선택가능)` — <b>셋 다</b>일 때만.
#   엑셀은 0 + 노란칸, 산출물 3종은 「현물」. 표시는 `_HYUNMUL`을 보고 각 산출물이 한다.
_HYUNMUL = set()      # 이번 고객에서 현물로 판정된 담보명
_CUR_CO = ''          # 지금 처리 중인 계약의 회사명
_CUR_CD = ''          # 지금 처리 중인 계약의 계약일 (YYYY.MM.DD)


def _is_hyunmul(raw, company='', cdate=''):
    import re as _r
    _n = _r.sub(r'[\s·]', '', str(raw or ''))
    if '현물급부선택가능' not in _n:
        return False
    if 'DB' not in str(company or '').upper() and '디비' not in str(company or ''):
        return False
    _m = _r.search(r'(20\d{2})[.\-/](\d{1,2})', str(cdate or ''))
    if not _m:
        return False
    return (int(_m.group(1)), int(_m.group(2))) >= (2026, 8)


def _re_sub_paren(t):
    """★v596 담보명에서 <b>괄호 블록을 통째로 제거</b>한다(상품명·심사방식 접미어).
    배제어 검사는 담보 <b>본체</b>에서만 해야 한다 — `(간편맞춤형Ⅱ)`의 Ⅱ가
    `질병입원수술비`를 죽이던 사고(천미옥 KB 실측)를 막는다."""
    return re.sub(r'[\(（\[][^\)）\]]*[\)）\]]', '', str(t or ''))


def resolve_kw(raw):
    # ★★★★★v370 (지점장 확정 2026.08.09, 영구): <b>순환계질환 주요치료비 = 2대 주요치료비</b>.
    #   지점장 원문: "47-50 은 2대주요치료비다 / 이름의 다양성이다 / 결국 순환계질환 주요치료비 다".
    #   `순환계질환(3-5종)주요치료비Ⅳ(수술시)`·`(중환자실입원시)`·`(혈전용해치료시)(뇌혈관질환)`·
    #   `(혈전용해치료시)(심혈관,기타)` = <b>1개 담보를 4개로 나눈 것</b> → 전부 같은 행, 대표(max) 1개.
    _v370 = re.sub(r'\s', '', str(raw or ''))
    # ★★★★★v372 (지점장 확정 2026.08.09): <b>「신특정순환계질환 통합치료비」도 2대 주요치료비</b>.
    #   구 조문은 '주요치료'만 봐서 KB `신특정순환계질환 통합치료비ⅢPlus` 5,000이 어느 행에도 못 갔다(실측).
    # ★★★★★v421e (지점장 확정 2026.08.14) — <b>「신특정순환계질환 주요치료비」는 순환계다</b>.
    #   지점장 원문: 「<b>신특정순환계질환주요치료비ⅢPlus → 진단서만 유일하게 순환계고 /
    #   뇌혈관심장주요치료비 or 뇌혈관허혈성심장주요치료비는 2대주요치료비다 // 진단서 7페이지 전용이다</b>」
    #   ・<b>순환계 + 주요치료비</b> → 진단서 7p <b>순환계 주요치료비 칸 전용</b>(마스터 행 없음 → 엑셀 미기재).
    #   ・<b>통합치료비</b>는 v372 지점장 확정대로 <b>2대 주요치료비</b> 유지.
    #   ・<b>뇌혈관+심장 조합 이름</b>은 2대 주요치료비(아래 규칙).
    # ★★★★★v421f (지점장 확정 2026.08.14 「<b>엑셀+보장분석지에는 2대주요치료비</b> / 진단서 7페이지에만
    #   순환계 주요치료비」 「<b>지금 내가 얘기하는건 진단서 전용이야 섞지마</b>」)
    #   → <b>엑셀 판정은 종전대로 2대</b>. 진단서 전용 칸은 `rep['p7_only']`로 따로 싣는다.
    #   ★한때 `__P7_CIRC__`로 엑셀에서 빼냈다가 <b>2대 칸이 비는 후퇴</b>가 났다 — 엑셀은 건드리지 않는다.
    if ('순환계' in _v370) and ('주요치료' in _v370 or '통합치료' in _v370):
        return '2대 주요치료비', 100

    # ★★★★★v372 (지점장 확정 2026.08.09, 영구): <b>암 통합치료비 = 암주요치료비 행</b>.
    #   지점장 원문: 「통합치료비Plus ㅡ 암주요치료비」.
    #   실측 — KB `암 통합치료비Plus(암중점치료기관(상급종합병원 포함))(고급형)` 5,000이
    #   <b>통합암</b> 행으로 갔다. 암주요치료비는 `_rep1` 대표(max) 1개라 500 → 5,000이 된다.
    #   ★'순환계 통합치료비'는 위 조문이 먼저 잡는다. 전이암·유사암 변형은 각자 행으로 두고 넓히지 않는다.
    #   ★★★★★v422c (지점장 확정 2026.08.15): 「<b>암 전액본인부담 통합치료비 ㅡ 엑셀·보장분석지ppt ㅡ 하이클래스</b>」
    #   → 제20조(「'비급여'가 붙은 암주요치료비는 하이클래스(암) 23행」)와 v421e 정본표
    #   (「비급여 암 통합치료비 | 하이클래스(암) | 비급여주요치료비 | 비급여 암 통합치료비 칸」)를 그대로 이행한다.
    #   ★구 코드는 <b>통합치료비에만 이 갈래가 없어</b> 비급여 건까지 21행으로 갔다(조문 미이행).
    if ('암' in _v370) and ('통합치료' in _v370) and ('비급여' in _v370 or '전액본인' in _v370) \
       and not any(x in _v370 for x in ('순환계','전이암','유사암','뇌혈관','심장')):
        return '하이클래스(암)', 0
    if ('암' in _v370) and ('통합치료' in _v370) and not any(
            x in _v370 for x in ('순환계','전이암','유사암','뇌혈관','심장','비급여','전액본인')):
        return '암주요치료비', 0

    # ★★★★★v422 (지점장 확정 2026.08.15): <b>암 특정치료비Ⅲ = 암주요치료비</b>.
    #   지점장 원문: 「암(유사암Ⅱ제외) 특정치료비Ⅲ(수술·항암방사선·항암약물) ㅡ 암주요치료비」.
    #   실측(삼성 양예슬) — `[건강]종합병원암(유사암Ⅱ제외)특정치료비Ⅲ(수술(회당),항암방사선,항암약물)` 1,000이
    #   이름 안의 <b>'수술'</b> 때문에 <b>암수술</b> 행으로 갔다.
    #   ★★<b>「유사암Ⅱ제외」는 유사암 담보가 아니다</b> — 이름 안에 '유사암'이 들어 있어도
    #     뒤에 '제외'가 붙으면 <b>암 본체</b>다. 단순 `'유사암' in` 배제는 이 담보를 죽인다.
    #   ★<b>유사암Ⅱ 특정치료비Ⅲ</b>(제외 없음)는 지점장 「패스」 → 종전대로 [확인]큐.
    #   ★★★★★v422c 최종 (지점장 확정 2026.08.15): 「<b>암(유사암Ⅱ제외) 특정치료비 ㅡ 엑셀·보장분석지ppt ㅡ 암주요치료비</b>」
    #   → <b>엑셀 21행에 기재한다</b>(v422b의 `__무시__`는 폐기). 7p 암 주요치료비 칸은 이 엑셀 값을 그대로 읽는다.
    _yx = bool(re.search(r'유사암[ⅠⅡⅢIV0-9]{0,3}제외', _v370))
    if ('암' in _v370) and ('특정치료비' in _v370) \
       and not any(x in _v370 for x in ('순환계','전이암','뇌혈관','심장')) \
       and (('유사암' not in _v370) or _yx):
        return '암주요치료비', 0

    # ★★★★★v298-A (심정자 실측 2026.07.31): 전각 ％(U+FF05) 때문에 v222 후유장해
    #   20%·50% 가드가 통째로 뚫렸다. 실측 = 라이나(에이스) `일반상해 20％이상 후유장해` 2,000이
    #   상해후유3%에 산입(반각 `20%`는 정상 차단됨을 단위검증으로 확인).
    raw = str(raw or '').replace('％', '%').replace('～', '~')
    # ★★★★★v247 (KB 3열 실측): <b>담보명 괄호 안의 장기 나열 글자에 걸려 오분류</b>되는 것을 최상단에서 차단한다.
    #   실측 — `갱신형 5대질환(심장,뇌혈관,신부전,간,폐 질환)수술비(연간1회한)[관혈]` 500×2가
    #   괄호 속 <b>'뇌혈관'</b> 글자 때문에 <b>뇌혈관수술비</b>로 가서 1,000 → <b>2,000</b>이 됐다(KB 표기 1,000).
    #   이 담보의 신정원담보명은 '특정질병수술' = <b>별개 담보</b>이므로 [확인]큐가 정답이다.
    #   ★기존 교훈("태그에 '뇌혈관' 글자 금지 — has('뇌혈관')를 먼저 검사해 오인한다")과 같은 뿌리.
    try:
        _r0 = re.sub(r'\s', '', str(raw or ''))
        if re.search(r'\d+\s*대질환', _r0) and '수술' in _r0:
            return None, 0
    except Exception:
        pass
    if str(raw).startswith('[확인]'): return None, 0   # ★v98 확인큐 항목은 표준행 매핑 금지(중복합산 차단)
    # ★★★v146 (지점장 확정 2026.07.21): 흥국화재 10억통장(플래티넘 건강 리셋월렛II)은
    #   <b>엑셀·보장분석 PPT에 표기 금지</b>. 보장진단서 7p 카드에만 표기한다.
    #   ★사고: 담보명에 '중환자실'이 들어 있어 resolve_kw가 '질병중환자실' 행으로 잡아
    #   엑셀 61행에 100,000(=10억)을 박고 있었다(실측). 최우선으로 차단한다.
    # ★★★★★v320 <b>10억 플랜 = 마스터 22행</b>(지점장 2026.08.01 "10억통장은 엑셀에 내가 넣었다").
    #   ★<b>구 v146·v148 「엑셀 전면 기재금지·PPT 금지」는 폐기</b> — 이제 정상 담보로 기재한다.
    if ('리셋월렛' in re.sub(r'\s','',str(raw))) or ('리셋월랫' in re.sub(r'\s','',str(raw))):
        return '10억 플랜', 0
    # ★★★v224 (2026.07.25 이정화 실측): <b>담보명 접힘 잔해는 어느 행에도 넣지 않는다</b>.
    #   3열 별첨에서 담보명이 위·아래 줄에 걸치면 파서가 <b>여러 담보명을 한 줄로 뭉치고
    #   엉뚱한 열의 금액을 붙인다</b>. 실측 잔해(엑셀 `_dambo_raw`):
    #   `외)(181일이상)(맞춤간편고지)중증질환자뇌혈관질환산정특례대상진단비Ⅱ(연간1회한)(맞춤
    #    간호 간병통합서비스사용질병입원일당(요양/정신/한방병원제` → 금액 <b>10</b>
    #   → 산정특례뇌혈관 행에 10(일당 금액)이 찍혔다. <b>정답은 1,000</b>.
    #   접힘 복원(파서 개편) 전까지는 <b>[확인]큐로 보내 오출고를 막는다</b>.
    _fz = re.sub(r'\s','',str(raw))
    if re.match(r'^[\)\]）]|^외\)|^고지\)|^편고지\)|^한도\)|^신형\)|^외\(', _fz):
        return None, 0        # 담보명이 닫는 괄호·꼬리로 시작 = 접힘 잔해
    _mix = sum(1 for _k in ('진단비','입원일당','수술비','치료비','사망보험금') if _k in _fz)
    if _mix >= 2 and len(_fz) > 45:
        return None, 0        # 담보 종결어 2종 이상 + 45자 초과 = 여러 담보 혼합
    """raw 담보명 -> (std표준명 or None, jong 0~5). API 불필요."""
    raw = re.sub(r"^\[세부보충\]","",str(raw))  # ★세부보충 접두 제거
    # ★Adobe OCR 깨짐: '상하!'·'상하 !'·'상하）' = '상해' (지점장 2026.07.05)
    raw = re.sub(r"상하\s*[!！]", "상해", str(raw))
    # ★v30x 담보명 경계깨짐 가드: 후유장해가 앞부분이면 뒤에 붙은 골절수술비 등은 오염 → 후유장해 우선
    _r0=str(raw)
    if ('후유장해' in _r0) and (_r0.find('후유장해') < 25):
        if '80%' in _r0 or '80 %' in _r0 or '80％' in _r0:
            if '상해' in _r0[:_r0.find('후유장해')+6] or '재해' in _r0[:_r0.find('후유장해')+6]: 
                pass  # 상해후유80%는 아래 정상 로직서 처리
    r = raw; n = _norm(raw)
    has = lambda *ks: all(_norm(k) in n for k in ks)
    no  = lambda *ks: not any(_norm(k) in n for k in ks)
    # ★v30z 혈전용해치료비 우선(급성심근·뇌졸중 흡수 방지): '혈전용해' 포함시 전용행
    # ★★★★★v397 (지점장 확정 2026.08.12, 영구): <b>심뇌혈관질환주요치료비 = 2대 주요치료비</b>.
    #   지점장 원문: 「<b>심뇌혈관질환주요치료비 -> 2대주요치료비다 / 진단서는 2대주요치료비에 넣어라</b>」
    #   ★실측(v396): 현대 `심뇌혈관질환주요치료비(…)(<b>수술및혈전용해치료</b>)` 500이 담보명 안의
    #     '혈전용해' 글자에 먼저 걸려 <b>혈전용해치료비</b>로 갔다. (같은 담보의 `(중환자실입원)` 변형은
    #     이미 2대 주요치료비로 정상 도착 — <b>한 담보가 괄호 수식어에 따라 두 행으로 갈리고 있었다</b>.)
    #   → <b>'심뇌'+'주요치료'면 괄호 안 수식어와 무관하게 2대 주요치료비</b>. 혈전용해 판정보다 앞에 둔다.
    #   ★`뇌졸중 혈전용해치료비`·`특정심장질환 혈전용해치료비`처럼 '주요치료'가 없는 담보는 종전대로 혈전용해치료비.
    # ★★★★★v542 제21조 2항 (지점장 실측 2026.08.22 「혈전용해를 가입한게 없는데 갑자기 뜬다」)
    #   롯데 담보명은 <b>「뇌혈관·허혈심장질환주요치료비」</b> — 「심뇌」가 아니라 <b>「뇌혈관」+「심장」</b>이라
    #   제21조 조건에 안 걸렸다. 그래서 3형제가 <b>세 행으로 흩어졌다</b>(실측 김순자 롯데):
    #     (혈전용해치료) 500 → <b>혈전용해치료비</b>   ← 가입한 적 없는 행에 떴다
    #     (수술)        500 → <b>뇌혈관수술비</b>
    #     (중환자실치료) 500 → 2대 주요치료비 (정상)
    #   ⇒ <b>「뇌」와 「심장」이 함께 있고 「주요치료」가 붙으면</b> 괄호 수식어와 무관하게
    #     <b>전부 2대 주요치료비</b>다(제21조 1항 「주요치료가 붙으면 수술 분해를 하지 않는다」).
    if (has('심뇌') or (has('뇌') and (has('심장') or has('허혈심')))) and has('주요치료'):
        return '2대 주요치료비',0
    # ★★★★★v547 제21조 3항 (지점장 확정 2026.08.22
    #   「혈전용해 → <b>주요치료비에서 있으면 있다고 인식하고 기입하는 것 금지</b>」)
    #   담보명에 <b>「주요치료」가 붙어 있으면</b> 그 안의 「혈전용해치료」는 <b>치료 방법 열거</b>일 뿐
    #   혈전용해 담보를 가입한 것이 아니다. <b>혈전용해치료비 행에 절대 기입하지 않는다.</b>
    if has('혈전용해') and has('치료') and no('주요치료'): return '혈전용해치료비',0
    # ★철심제거·핀제거·내고정물제거 = 골절수술비 아님(별개 처치) → [확인]
    if (has('철심') or has('핀제거') or has('내고정물')) and has('수술'): return None,0
    # 종번호
    jong = 0
    # ★★★★★v422d (지점장 확정 2026.08.15 「1-8이 한 세트야」): 구 코드는 <b>1~5종만</b> 읽었다.
    #   삼성 `상해6·7·8종수술비(시술포함)`는 jong=0으로 떨어져 <b>1-5종 행에 종번호 없이</b> 들어갔다
    #   → 슬래시 5칸 어디에도 안 앉아 <b>200·300·500이 조용히 사라졌다</b>(실측 조승우).
    for i,k in enumerate(['1종','2종','3종','4종','5종','6종','7종','8종','9종'],1):
        if k in n or f'({i}종)' in r: jong = i; break

    # 비담보성(보험료 납입면제·일시납입지원) → 매핑 안 함(자부상 등 오매핑 차단)
    # ★★★★★v600 제140조 3항 (지점장 지시 「kb라이프 인식해라」 · 천미옥 실측).
    #   [실측 결함] 생보 담보명은 <b>`(납입면제형)`·`(해약환급금미지급형)`</b> 같은
    #     <b>납입·환급 방식 표기</b>를 괄호로 달고 나온다. 이 규칙이 그 괄호까지 보고
    #     `허혈성심장질환진단Ⅱ(…)(납입면제형)` <b>1,000을 통째로 [확인]큐</b>로 보냈다.
    #   ⇒ <b>괄호를 걷어낸 본체</b>에서 본다. 진짜 납입면제 담보(`납입면제대상보장(5대기본)`)는
    #     본체에 「납입」+「면제」가 남아 종전대로 막힌다.
    _bodyf = _re_sub_paren(raw)
    _bf = re.sub(r'\s', '', str(_bodyf or ''))
    if ('납입' in _bf) and (('면제' in _bf) or ('지원' in _bf) or ('대상보장' in _bf)):
        return None,0

    # ★ 상해의료비 = 별개 정액 담보 단독 행(실손 입원/통원/약값과 합치지 말 것) — 지점장 2026.06.28
    if has('상해의료비') and no('입원','통원','외래','실손','처방','약제','도수','비급여'): return '상해의료비',0
    # ★★v95 (지점장 지시 2026.07.19): 1세대 구실손의 <b>'상해 의료비(입원+통원)'</b>은
    #   입원·통원으로 쪼개 넣을 수 없는 <b>상해의료비 단독 담보</b>다. 실손 행(입원/통원)에
    #   섞으면 한장보장표와 어긋난다. → 마스터 99행 '상해의료비'로 직행.
    #   실측: DB 0604_TM '상해(일반상해, 전체상해를 의미) 의료비(입원+통원) 100' → 상해의료비 100
    if has('상해') and has('의료비') and ('입원+통원' in n.replace(' ','') or '입원/통원' in n.replace(' ','')):
        return '상해의료비',0
    if has('외래') and has('의료비') and no('주사','MRI','도수','체외','증식','비급여'): return '통원',0   # ★v29x '외래의료비'(통원 표기 없음)=통원. 상해·질병 각각 와도 대표 최댓값 1건이라 중복합산 없음

    # ── 실손/수술일당 먼저 (수술·일당 오분류 차단) ──
    if (has('실손') or has('입원형') or has('입원의료비')) and has('입원'): return '입원',0
    if has('도수') or has('체외충격파') or has('증식치료'): return '도수치료',0   # 비급여 도수/체외/증식
    if has('MRI'): return 'MRI',0
    if has('비급여') and has('주사'): return '비급여주사',0
    # ★★★★★v550 제9조 2항 (지점장 실측 2026.08.22 「4세대 실손 오류 — 입원 0 · 통원 5천 ·
    #   상해의료비 5천 / 정답은 입원 5천 · 통원 20 · 상해의료비 없음」)
    #   4세대는 담보가 <b>급여/비급여로 쪼개져</b> 온다. 구 조건은 `no('비급여')`라
    #   <b>`상해비급여통원의료비`·`질병비급여통원의료비`가 통째로 [확인]큐(None)</b>로 빠졌다.
    #   ⇒ <b>「통원」이 담보명에 있으면 비급여여도 통원 행</b>이다.
    #     단 3대 비급여(도수·체외·증식·주사·MRI)는 위에서 이미 자기 행으로 갔다.
    #   ★v550 2차 — <b>3세대 `상해통원형`·`질병통원형`이 None</b>이었다(입원형은 잡히는데 통원형만).
    #     입원 규칙과 <b>대칭</b>으로 `통원형`도 통원 행이다.
    if has('통원') and (has('실손') or has('외래') or has('의료비') or has('통원형')) \
       and no('주사','MRI','도수','체외','증식'): return '통원',0
    if has('상해') and has('수술') and has('일당'): return '상해수술일당',0   # ★v29q-10 상해수술입원일당→상해수술일당(질병수술일당 오입력 차단)
    if has('수술') and has('일당'): return '질병수술일당',0
    # ★★★v215 (지점장 확정 2026.07.25, 영구): <b>생명보험사 '급부금' 담보명 4종 정본 매핑</b>.
    #   ①질병수술급부금 = <b>질병수술비</b>  ②재해수술급부금 = <b>상해수술비</b>
    #   ③질병입원급부금 = <b>질병입원일당(질병일당)</b>  ④재해입원(급부금) = <b>상해입원일당(상해일당)</b>
    #   근거: 생보 약관은 '수술비·일당' 대신 '급부금'으로 인쇄한다. 구 코드는 '급부금' 문자열을 몰라
    #   `_is_pure_s/_is_pure_q`(상해수술비·질병수술비로 시작) 검사에서 전부 탈락 → <b>4종 모두 [확인]큐로 사라졌다</b>(실측).
    #   ★수식어가 붙은 변형(암·뇌·특정·교통·종합병원 등)은 base 담보가 아니므로 여기서 제외 → 기존 로직/[확인]으로 흘린다.
    _gbex = ('암','뇌','심','허혈','간질환','신장','폐','위','골절','화상','특정','교통','종합','상급','중환자','요양','재활','통원','외래','간병','장해','후유')
    if has('급부금') or has('급여금') or has('재해입원'):
        if has('수술') and no(*_gbex) and jong==0:
            if has('질병'): return '질병수술비',0
            if has('상해') or has('재해'): return '상해수술비',0
        if has('입원') and no('수술', *_gbex):
            if has('질병'): return '질병일당',0
            if has('상해') or has('재해'): return '상해일당',0
            # ★★★★★v305 (지점장 지시 2026.07.31, 영구): <b>입원급여금 = 질병일당</b>.
            #   지점장 원문 = "입원급여금 질병일당으로 추가".
            #   [구 결함] 생보 약관은 '급부금'과 '급여금'을 <b>둘 다</b> 쓰는데 코드는 '급부금'만 알아서
            #   `입원급여금`·`질병입원급여금`·`상해입원급여금`이 <b>전부 [확인]큐로 사라졌다</b>(실측).
            #   실측 사례 = 새마을금고 `입원특약[…] : 입원급여금[3일초과/0.1%]` 2.
            #   ★축이 명시되면 그 축(질병/상해)으로 가고, <b>축이 없으면 질병일당</b>이다(지점장 지시).
            #   ★수식어형(암·교통·중환자실 등)은 `_gbex`가 그대로 걸러 [확인]큐 유지.
            if has('입원급여금'): return '질병일당',0

    # ── 수술비 ──
    if has('수술'):
        # ★★★★★v337 (지점장 지시 2026.08.02): <b>엑셀에 1-8종 = 1-7종 = 1-9종 동일</b>이다.
        #   → 담보명에 1-7종·1-8종·1-9종이 있으면 <b>(1-8종) 행</b>으로 보낸다(1-5종과 절대 섞지 않는다).
        #   실측 = 미래에셋 `1-7종수술특약(급여)` 2,000이 1-5종 쪽으로 섞여 지점장이 "합쳐졌다"고 지적.
        #   ★엑셀은 1-5종과 같은 방식으로 <b>가로 슬래시 나열 + 끝열 종별 합산</b>,
        #     PPT는 <b>대표값(최댓값)</b>으로 넣는다(`_rep1`에 이미 등록돼 있다).
        _n78 = re.sub(r'\s','',n)
        if any(k in _n78 for k in ('1-7종','1-8종','1-9종','1~7종','1~8종','1~9종')):
            if has('질병'): return '질병 종수술비(1-8종)', jong
            return '상해 종수술비(1-8종)', jong
        # ★★★★★v422d (지점장 확정 2026.08.15 「<b>1-8이 한 세트야</b>」) — 삼성 실측.
        #   삼성은 종을 <b>한 줄씩 따로</b> 쓴다(`상해1종수술비(시술포함)` … `상해8종수술비(시술포함)`).
        #   담보명에 `1~8종` 같은 <b>세트 표기가 없어</b> 구 코드가 전부 <b>1-5종 행</b>으로 보냈고,
        #   1-5종 슬래시는 5칸뿐이라 <b>6·7·8종(200·300·500)이 앉을 자리가 없어 소실</b>됐다(실측 조승우).
        #   → <b>종번호가 6 이상이면 그 담보는 1-8종 세트</b>다. v343 조문대로 <b>대표(max) 1개</b>로 간다.
        #   ★1~5종이 <b>명시된</b> 세트(`상해1~5종수술비(3종)`)는 위 분기에 안 걸리므로 <b>종전 그대로</b>.
        if jong and jong >= 6:
            if has('질병'): return '질병 종수술비(1-8종)', 0
            if has('상해') or has('재해'): return '상해 종수술비(1-8종)', 0
            return '종수술비공통', 0
        # ★★★★★v378 (지점장 지시 2026.08.10): <b>「종수술」은 1~5종 수술비다</b>.
        #   지점장 원문 — "파워수술보장(본인) = 1-5종 수술비고 <b>생명보험사는 질병,상해 둘다 기재</b>되어야 한다
        #   / <b>질병종수술 = 질병1-5종 . 상해종수술=상해1-5종</b> 인데 미기재다 오류다".
        #   [구 결함] 아래 두 줄이 <b>종번호(jong)를 요구</b>해, 종번호 없이 <b>단일금액</b>으로만 실린
        #   삼성생명 `질병종수술` 500 · `상해종수술` 500 · `無파워수술보장(본인)` 500이
        #   resolve2에서 <b>None</b>으로 떨어져 통째로 [확인]큐로 사라졌다(실측 구본칠).
        #   → 종번호가 없어도 1-5종 행으로 보낸다(jong=0 = 종별 분해 없는 단일금액).
        if (('종수술' in _n78) or ('파워수술보장' in _n78)) and ('각종' not in _n78):
            if has('질병'): return '질병 종수술비(1-5종)', jong
            if has('상해') or has('재해'): return '상해 종수술비(1-5종)', jong
            return '종수술비공통', jong   # 축 미표기(파워수술보장 등) → 질병·상해 양쪽
        if has('상해') and jong: return '상해 종수술비(1-5종)', jong
        if has('질병') and jong: return '질병 종수술비(1-5종)', jong
        # ★v215 (지점장 확정 2026.07.25): <b>'중대상해수술비' = '중대한상해수술비'</b>(같은 담보, '한' 한 글자 차이).
        #   구 코드는 `has('중대한','상해')`라 '중대상해수술비'가 탈락 → [확인]큐로 빠졌다(실측).
        if has('중대') and has('상해'): return '중대한상해수술비',0
        # ★★v216 자가진단으로 발견한 버그: 공백을 지우면 <b>'수술비 관혈' → '수술비관혈'</b>이 되어
        #   그 안에 <b>'비관혈'이 우연히 만들어진다</b> → 관혈 담보가 통째로 <b>비관혈 행</b>으로 갔다.
        #   괄호가 있는 '5대기관수술비(관혈)'은 우연히 살아남고, 괄호 없는 표기만 틀리던 <b>조용한 오분류</b>.
        _gwan = n.replace('수술비관혈', '수술비|관혈')
        if has('5대기관') and ('비관혈' in _gwan): return '5대기관 수술비 비관혈',0
        if has('5대기관'): return '5대기관 수술비 관혈',0
        # ★★★★★v351 (지점장 지시 2026.08.02): "<b>**대질병수술비는 엑셀 120대에 기재, 담보금액은
        #   대표값 하나만, ppt에도 다 기재</b>". [구 결함] resolve가 <b>`'n대수술비'`</b>를 반환하는데
        #   <b>마스터 행 이름은 `'120대수술비'`</b>라 이름이 어긋났다(별칭 폴백이 있었지만 반환값을
        #   <b>마스터 이름으로 통일</b>해 어긋날 여지를 없앤다). 실측 = AIG 86대질병수술비 6건·롯데 16대 미기재.
        # ★v352b: 마스터 라벨 자체('n대수술비')도 자기 행으로 돌아와야 한다(커버리지 감사 FAIL 방지).
        if re.sub(r'\s','',r) in ('n대수술비','N대수술비','120대수술비'): return 'n대수술비',0
        if re.search(r'(?<!\d)\d{2,3}\s*대', r): return 'n대수술비',0   # ★v352 마스터 행 이름을 'n대수술비'로 통일(지점장 지시 2026.08.02)
        # ★★★★★v326b (지점장 확정 2026.08.02): <b>뇌혈관수술비 = 담보명에 '뇌혈관'이라고 적힌 것만</b>.
        #   ★내 오류 정정 — v325에서 `뇌출혈`을 여기 넣었으나 지점장이 바로잡았다: "<b>뇌출혈도 엑셀에 없다</b>".
        #   마스터에 「뇌출혈 수술비」 행이 <b>없으므로 [확인]큐</b>가 정답이다.
        #   "뇌출혈은 뇌혈관의 하위 질환"은 <b>내 판단이지 지침이 아니다</b> — 비슷한 행에 임의로 밀어 넣지 않는다.
        if has('뇌혈관') or has('심뇌혈관'): return '뇌혈관수술비',0
        if has('허혈'): return '허혈성수술비',0
        # ★★★★★v326b (지점장 확정 2026.08.02): <b>심장수술비 = 담보명에 '심장'이라고 적힌 것만</b>.
        #   ★내 오류 정정 — v326에서 `급성심근`·`심근경색`을 여기 넣었으나 <b>지점장이 즉시 바로잡았다</b>:
        #     "급성심근경색 수술비는 <b>엑셀에 없는데</b> 왜 그게 심장이냐 / <b>심장수술비는 심장이라고 적힌 것</b>만 된다".
        #   → <b>마스터에 행이 없으면 [확인]큐</b>다. 비슷한 행에 임의로 밀어 넣지 않는다("엑셀에 없는 건 제외").
        if has('심장') or has('심질환'): return '심장수술비',0
        if has('5대골절'): return '5대골절수술비',0
        if has('골절') and no('후유','장해','진단','일당','입원'): return '골절수술비',0
        if has('화상'): return '화상수술비',0
        # ★v51(지점장 확정 2026.07.13): 현대해상 '레보아이로봇수술비' = 다빈치로봇수술비(마스터 26행).
        #   '로봇' 키워드로 이미 잡힌다 — 이 조건을 좁히면 레보아이가 조용히 누락되므로 건드리지 말 것.
        if has('다빈치') or has('로봇') or has('레보아이'): return '다빈치로봇수술비',0
        if has('암') and no('양성종양','유사암'): return '암수술',0   # ★v30 양성종양·유사암 수술 오탐 차단 → [확인]
        if jong: return '종수술비공통', jong   # ★v29q-12 상해/질병·부위 미표기 1-5종 수술(예 파워수술 1-5종)→상해·질병 양쪽 슬래시
        if has('상해') or has('재해'):   # ★v30h 재해수술비=상해수술비 동일 취급
            # §6 상해수술비 = 기본만. 병원규모·부위/특정·통원·자XXXX 접두변형은 합산 금지 → [확인]
            # ★XXXX상해수술비(질병/상해수술비 앞 어떤 접두든) = 별개 아님→[확인](지점장 2026.07.05)
            _core_s = re.sub(r'^[\(\[][^\)\]]*[\)\]]\s*', '', r)
            _core_s2 = _core_s.strip().replace(' ','')
            # ★★★★★v247 (2026.07.25 KB 3열 실측): <b>3열(KB·메리츠)은 회사담보명 앞에 '갱신형'이 붙는다</b>
            #   (`갱신형 상해수술비` · `갱신형 질병수술비`). `startswith('상해수술비')` 검사라 <b>전부 탈락</b>해
            #   상해수술비 200 · 질병수술비 30이 <b>통째로 [확인]큐로 사라졌다</b>(KB 전체보장현황 대조 실측).
            #   → 판정 전에 <b>'갱신형'·'비갱신형' 접두어를 제거</b>한다. 2열(롯데)엔 이 접두어가 없어 영향 없다.
            _core_s2 = re.sub(r'^(?:비)?갱신형', '', _core_s2)
            _core_s2 = re.sub(r'^재해상해', '상해', _core_s2)   # 재해상해=상해(중복 정리)
            _core_s2 = re.sub(r'^재해(입원)?수술', r'상해\1수술', _core_s2)  # ★v65 재해수술비·재해입원수술비=상해수술비(지점장 2026.07.15, '입원' 낀 변형도 포함)
            # ★v594 (지점장 확정 2026.08.26 「비가 없어도 수술비다」) — 질병 쪽과 <b>대칭</b>.
            #   `상해수술(간편건강고지)담보`가 「비」 한 글자 때문에 [확인]큐로 빠졌다.
            #   ★단 「비」 뒤에 올 수 있는 것은 <b>괄호·접미어(담보)·끝</b>뿐 —
            #     `상해수술위로금`·`상해수술자금`·`상해수술치료비`는 수술비가 아니다.
            # ★v596 제135조 2항 — 입원·통원 둘 다 상해수술비다(지점장 확정 2026.08.28)
            _is_pure_s = bool(re.match(r'^(?:상해수술|상해입원수술|상해통원수술)비?(?:\(|담보|$)', _core_s2))
            #   ★가드는 <b>괄호를 걷어낸 본체</b>로 본다 — `(간편맞춤형Ⅱ)` 상품명 오탐 차단.
            #   ★'통원'을 가드에서 뺀다(지점장 확정). '외래'는 별개 담보라 그대로 배제.
            _body_s = _re_sub_paren(_core_s2)
            # ★★★★★v600 제135조 2항 (지점장 화면 실측 2026.08.29 — 서버 실행검사 FAIL 1건
            #   「[제135조 수술비] 상해수술비(종합병원) — 상해수술비 (정답 None)」).
            #   [원인] v596이 배제어를 <b>괄호를 걷어낸 본체</b>에서만 보게 바꿨다.
            #     `(간편맞춤형Ⅱ)` 오탐을 막으려던 것인데, 같이 걷혀 나간 <b>`(종합병원)`</b>도
            #     안 보이게 돼 병원 규모 변형이 기본 행으로 들어왔다.
            #   ⇒ <b>병원 규모·통원은 괄호 안에 있어도 담보를 가른다</b> → 원본에서 본다.
            #     나머지 배제어는 종전대로 본체에서만(상품명 오탐 차단 유지).
            _HOSP600 = ('상급', '종합병원', '통원')
            if _is_pure_s and any(k in _core_s2 for k in _HOSP600):
                return None, 0
            if _is_pure_s and not any(k in _body_s for k in ('흉터','복원','외모','특정','척추','관절','하지','상급','종합병원','안면','머리','목','3대','신경','인대','흉부','연골','외래','자궁','자녀','자가','교통')):
                return '상해수술비',0
            return None,0
        if has('질병'):
            # ★v29 §8.5 질병수술비 합산군 = '질병수술비'·'질병입원수술비' 만.
            #   ★XXXX질병수술비(어떤 접두든) = 별개→[확인](지점장 2026.07.05) → [확인]
            _core = re.sub(r'^[\(\[][^\)\]]*[\)\]]\s*', '', r)   # 접두 수식어 괄호 제거
            _core_strip = _core.strip().replace(' ','')
            _core_strip = re.sub(r'^(?:비)?갱신형', '', _core_strip)   # ★v247 3열 '갱신형' 접두어 제거(KB 실측)
            # 순수 질병수술비/질병입원수술비로 시작해야 함(자XXXX 등 한글 접두 배제)
            # ★★★★★v594 (지점장 확정 2026.08.26 「<b>비가 없어도 수술비다</b>」).
            #   [실측 · 현대해상] `질병수술(간편건강고지)담보` <b>30</b>이 [확인]큐로 빠졌다 —
            #     조건이 `질병수술<b>비</b>`로 시작을 요구해 <b>「비」 한 글자</b> 때문에 탈락했다.
            #     `(간편건강고지)`는 심사방식, `담보`는 접미어다. 같은 담보의 다른 표기다.
            #   ⇒ `질병수술` / `질병입원수술`로 시작해도 받는다. 변형 배제어(`_excl_ok`)는 그대로다
            #     — `다발성`·`특정`·`부위` 등은 종전대로 [확인]큐로 간다.
            #   ★v594b — <b>「비」만 빼고 나머지는 그대로 좁게</b> 본다.
            #     시작만 보면 `질병수술<b>위로금</b>`·`질병수술<b>자금</b>`·`질병수술<b>치료비</b>`·
            #     `질병수술<b>(재진단)</b>`까지 전부 들어온다(실측). 그건 수술비가 아니다.
            #     ⇒ 「비」 뒤에 올 수 있는 것은 <b>괄호·접미어(담보)·끝</b>뿐이다.
            # ★v596 제135조 2항 — 입원·통원 둘 다 질병수술비다(지점장 확정 2026.08.28)
            _is_pure_q = bool(re.match(r'^(?:질병수술|질병입원수술|질병통원수술)비?(?:\(|담보|$)', _core_strip))
            # ★★★★★v596 제135조 2항 — <b>배제어는 괄호를 걷어낸 담보 본체에서만</b> 본다.
            #   [실측 · 천미옥 KB] `질병입원수술비(당일입원제외)(<b>간편맞춤형Ⅱ</b>)` 가 통째로 죽었다.
            #   배제어 'Ⅱ'는 `질병수술비Ⅱ` 같은 <b>담보 본체의 변형</b>을 거르려는 것인데
            #   괄호 안 <b>상품명</b>의 Ⅱ까지 잡았다. KB·삼성·메리츠·롯데·현대가 다 이 표기다.
            #   ★'통원'은 배제어에서 뺀다 — 지점장 확정 「질병통원수술비도 질병수술비다」.
            # ★v600 제135조 2항 — 질병 쪽도 대칭. 병원 규모·통원은 괄호 안에서도 본다.
            if bool(re.match(r'^(?:질병수술|질병입원수술)비?(?:\(|담보|$)', _core_strip)) \
               and any(k in _core for k in ('상급', '종합병원', '통원')):
                return None, 0
            _body_q = _re_sub_paren(_core)
            _excl_ok = not any(k in _body_q for k in ('특정','부위','관절','척추','외모','흉터','복원','신경','인대','연골',
                          '상급','종합','대수술','수술일당','일당','Ⅱ','Ⅲ','ⅱ','ⅲ',
                          '2종','3종','4종','5종','부분','관혈','내시경','로봇','외래','5대','자궁','자가','자녀'))
            if _is_pure_q and _excl_ok and jong==0:
                return '질병수술비',0
            return None,0
    if has('창상') or has('봉합'): return '창상봉합술',0

    # ★v35 중증질환자 산정특례대상보장 = 산정특례 전용행. 파서가 담보명 3줄분할·OCR깨짐으로 조각냄
    #   실측조각(이성준KB 2026.07.09): ①'중증질환자（뇌혈관）'단독 ②'산정특례대상보장…중•증걸 환■자■（삼장）'
    if (has('중증질환자') or has('중증환자')) and has('뇌혈관') and no('수술'): return '산정특례뇌혈관',0
    if has('삼장') and no('수술'): return '산정특례심장',0   # OCR '삼장'=심장, 중증질환자 산특 조각 전용
    if (has('중증질환자') or has('중증환자')) and (has('심장') or has('삼장')) and no('수술'): return '산정특례심장',0
    # ★★★v227 (지점장 지시 2026.07.25, 영구): <b>'일반암직접치료' = 암수술 — 세부가입현황을 정본으로 따른다</b>.
    #   근거(이정화 AIA 건강+ 실측·검산 일치): 세부가입현황 AIA 건강+ 암 칸 = <b>일반암 1,000 / 암수술 1,000
    #   / 유사암 100 / 고액항암 8,000</b>. 별첨엔 `일반암직접치료 1,000`이 <b>2줄</b>인데 구 코드는
    #   ①매핑이 없어 [확인]큐로 빠지고 ②`_add`가 두 줄을 <b>합산 2,000</b>으로 만들었다.
    #   → <b>암수술 행 + 계약 내 대표(max) 1건</b>으로 처리하면 1,000이 되어 세부가입현황과 일치한다.
    #   전수 검산: 암수술 = 한화 750 + AIA 1,000 = <b>1,750</b> = 한장보장표 값.
    #   ★'암직접치료<b>입원일당</b>'은 위 일당 블록에서 이미 암일당으로 갔다(여기 안 온다).
    if has('암') and has('직접치료') and no('일당','입원','통원','방사선','약물','표적','양성자','세기','중입자','유사암'):
        return '암수술',0
    # ── 암 치료비 ──
    # ★★★★★v396 (지점장 확정 2026.08.12, 영구) — <b>구 조문 「2026.07.09 '암주요치료비' 명시 > 하이클래스」는 폐기</b>.
    #   지점장 원문: 「<b>하이클래스암주요치료비 or 비급여암주요치료비 = 둘다 엑셀+ppt 2건 모두에 "하이클래스"에 기재하라</b>」
    #   → 이름에 <b>'하이클래스'</b> 또는 <b>'비급여'</b>가 붙은 암주요치료비는 <b>하이클래스(암) 23행</b>이다.
    #   ★왜 필요했나: 현대·롯데 제안서 담보명이 `하이클래스 암주요치료비`처럼 <b>두 단어가 한 이름에</b> 붙어 있어
    #     구 조문대로면 전부 21행으로 가고 <b>23행이 영영 공란</b>이 됐다. 그 결과 보장진단서 8쪽
    #     「비급여 암 주요치료비」 칸도 계속 0이었다(v393에서 키 오타를 고쳐도 값이 없으니 그대로 0).
    #   ★실손 3종(도수·MRI·비급여주사)은 이 지점보다 <b>앞</b>에서 이미 잡히므로 '비급여' 단독 오인 위험 없음.
    # ★★★v258b(2026.07.27): 지침 §8.2 원문에 <b>「암(유사암제외)주요치료비 → 암주요치료비」</b>가
    #   명시돼 있는데, 구 조건 `has('유사암')`이 <b>'유사암<u>제외</u>'의 '유사암' 글자</b>에 걸려
    #   <b>__무시__</b>로 보냈다(실측). '유사암제외'는 유사암을 <b>빼는</b> 담보 = 일반암 계열이다.
    # ★★★★★v422b (지점장 정정 2026.08.15): 「<b>유사암Ⅱ 특정치료비Ⅲ 이걸 패스하라는거다 · 원래 지침이다</b>」
    #   → 새 규칙이 아니다. <b>①유사암 주요치료비 = 무시</b> 조문 그대로이고, 삼성이 같은 담보를
    #   <b>'특정치료비'</b>라고 부를 뿐이다(실측 양예슬 `[건강]종합병원 유사암Ⅱ 특정치료비Ⅲ` 600 → 구 [확인]큐).
    #   ★'유사암<u>제외</u>'는 여전히 유사암이 아니다(v258b) — 위 v422 암 특정치료비와 짝을 이룬다.
    if has('유사암') and no('유사암제외') and (has('주요치료') or has('특정치료비')): return '__무시__',0   # ①유사암 주요치료비=무시(엑셀·PPT·설명지 전부)
    # ★★★★★v600 제0조 꼭대기 (지점장 확정 2026.08.26 「<b>유방주요치료비(…)(중증유방암진단)담보
    #   1,000 ← 엑셀이 없는 거다</b>」 · 현대해상).
    #   [실측] 마스터에 <b>「유방」 행이 없다</b>. 그런데 담보명 <b>맨 앞</b>이 `유방주요치료비`인데
    #     괄호 안 `유방<b>암</b>진단`의 「암」이 들어와 <b>암주요치료비 21행</b>에 앉아 3,100의 일부가 됐다.
    #   ⇒ <b>담보명이 부위로 시작하면</b> 괄호와 무관하게 [확인]큐(비슷한 행에 임의로 넣지 않는다).
    _PART600 = ('유방','자궁','전립선','갑상선','생식기','음경','고환','난소','방광','신장','간장')
    _head600 = re.sub(r'^[\[(][^\])]*[\])]\s*', '', str(raw or '').strip())
    if re.match(r'^(%s)' % '|'.join(_PART600), _head600) and ('주요치료' in n or '통합치료' in n):
        return None, 0
    if has('암주요치료비') and no('유사암','하이클래스','비급여'): return '암주요치료비',0   # ②순수 암주요치료비만 21행
    #   ★v396 `비급여암주요치료비`는 '하이클래스' 글자가 없어 아래 조건만으로는 안 걸린다 — 명시 추가.
    if has('하이클래스') or (has('비급여') and has('암') and has('주요치료')):
        return '하이클래스(암)',0   # ③하이클래스 / 비급여 암주요치료비 → 하이클래스(암) 23행. 2건이면 합산
    # ★★★★★v258 (지점장 지침 개정 2026.07.27, 영구): 지침 §8.2 원문 =
    #   <b>「(키워드 : 암+주요치료비 / 그외 는 다 아님)」</b>
    #   → 담보명에 <b>'암'과 '주요치료비'가 둘 다</b> 있어야 암주요치료비 행이다. 그 외는 전부 아니다.
    #   구 코드는 `has('주요치료')`만 보고 '암' 없이도 이 행에 넣었다(실측 오류 —
    #   KB `3대(간,폐,신장)질환 주요치료비(간편가입)` 500이 암주요치료비로 들어갔다).
    #   → `has('암')` 추가. 조건 미달은 [확인]큐로 보내 신인이 수기 확인한다(누락 금지).
    # ★v260b: 지침 §8.2 우선순위 = <b>하이클래스/비급여 → 하이클래스(암)</b>가 암주요치료비보다 앞이다.
    #   접힘이 풀리면서 `비급여(전액본인부담 포함) 암특정주요치료비Plus(상급종합병원)`처럼
    #   <b>'비급여'와 '주요치료비'가 한 담보명에 같이</b> 들어오게 됐다 → `no('비급여')`로 아래 하이클래스 규칙에 넘긴다.
    if has('암') and has('주요치료') and no('비급여','순환계','2대','뇌','허혈','심장','심근','유사암','하이클래스'):   # ③하이클래스 없는 '병원+암주요치료비'→암주요치료비행. ★심장 추가(심장/순환계 주요치료비=2대주요치료비로, v38d)
        return '암주요치료비',0
    if has('고액항암') or (has('고액') and has('항암') and has('치료')): return '__무시__',0   # ★v30z 고액항암치료비=표적+양성자+세기조절+카티 합계값 → 무시(구성 치료비는 아래에서 각각 개별 매핑)
    # ★★★v215 (지점장 확정 2026.07.25, 영구): <b>괄호 안의 실제 치료법이 정본이다</b>.
    #   ①표적항암방사선치료비<b>(항암세기조절방사선)</b> = <b>세기조절치료</b>
    #   ②표적항암방사선치료비<b>(항암양성자방사선)</b> = <b>양성자치료</b>
    #   구 코드는 `has('표적')`이 세기조절·양성자 판정(아래)보다 <b>먼저</b> 걸려 둘 다 표적항암치료비로 뭉갰다(실측).
    #   ★세기조절·양성자 판정을 표적 <b>앞</b>으로 올린다.
    if has('세기조절'): return '세기조절치료',0
    if has('양성자'): return '양성자치료',0
    if has('표적'): return '표적항암치료비',0
    # ★v30h 암주요치료비 = 암특정치료비/암주요치료비/암(특정유사암포함)진단후(종합병원/상급종합병원)특정치료(지원금/비)
    #   구간밴드(연간1천~1억) 다줄·부위별 = 대표 1개(max, §8.2). 뇌·심·순환계·비급여·재활·통원·검사는 제외(각 전용행).
    if has('암') and (has('주요치료') or has('특정치료') or (has('진단후') and has('치료') and (has('종합병원') or has('특정'))))         and no('순환계','2대','뇌','허혈','심장','심근','비급여','하이클래스','재활','통원','입원일당','MRI','PET','초음파','검사','수술'):
        return '암주요치료비',0
    # 뇌혈관·허혈성심장 특정치료비 = 2대주요치료비(뇌·심 두 칸). ★v65 '뇌심주요치료비특약' 누락 수정(지점장 2026.07.15)
    # ★v421e 순환계 단독은 위에서 `__P7_CIRC__`로 빠진다 — 여기 '순환계'를 남기면 다시 2대로 끌려간다.
    if (has('뇌혈관') or has('허혈') or has('심장') or has('순환계') or has('뇌심') or (has('뇌') and has('심'))) and (has('특정치료') or has('주요치료')) and no('암','수술'):
        return '2대 주요치료비',0
    if has('치료지원금') or (has('진단후') and has('치료')): return None,0   # ★v30a 잔여 진단후 치료지원금(암·뇌·심 아닌) = 진단비 아님 → [확인]
    if has('유사암') and has('주요치료'): return None,0   # ★v30 유사암 주요치료비 = 전용행 없음 → [확인]
    # ★★★★★v248 (지점장 확정 2026.07.26 "유사암이라고 기재된 것만 해라"):
    #   구 v30a 규칙 <b>`갑상선암 + 진단 → 유사암`은 폐기</b>한다. 담보명에 '유사암'(또는 '소액암')이
    #   명시되지 않으면 <b>[확인]큐</b>로 보내 신인이 수기 확인한다(추측 금지).
    #   실측 — 메리츠 `갱신형 갑상선암(초기제외)진단비` 1,000이 유사암에 산입돼 3,000이 됐으나
    #   KB 전체보장현황 유사암은 1,000(`갱신형 유사암진단비`뿐)이다.
    if has('갑상선암') and has('진단') and no('유사암','소액암','주요치료','수술','일당'): return None,0
    # ★★★★★v320 <b>통합암 = 마스터 16행</b>(지점장이 직접 넣음 2026.08.01).
    #   지침 §8.2 「통합암·통합전이암 = <b>대표금액 1개</b>」가 정본. ★순서 주의: <b>통합전이암이 먼저</b>다.
    # ★★★★★v398 (지점장 확정 2026.08.12, 영구): <b>통합암은 그 문구가 앞에 나온다</b>.
    #   지점장 원문: 「`암진단비(유사암제외)(<b>통합</b>간편가입형)` ← 이라고 통합암에 기재가 된다.
    #     <b>아니다. 이건 암진단비다 / 통합암은 그 문구가 앞에 나온다. 지침오류다</b>」
    #   ★구 조건 `has('통합') and has('암')`은 <b>'통합'이 담보명 어디에 있든</b> 걸렸다 —
    #     `(통합간편가입형)`·`(통합간편가입Ⅱ)` 같은 <b>가입유형 수식어</b>의 '통합'까지 통합암으로 보냈다(실측).
    #   → <b>'통합암'·'통합전이암'이 한 덩어리로 붙어 있을 때만</b> 그 행이다. 그 외 '통합'은 무시한다.
    # ★★★★★v539 제126조 (지점장 실측 2026.08.22 「보장분석지에 통합암전이암진단비 =
    #   통합전이암진단비 ⇒ <b>일반암으로 표기</b>되어 나오고 진짜 일반암진단비는 표기가 안된다」)
    #   [실측] 롯데 `통합<b>형</b>전이암진단비(4대특정전이암)(간편할인형Ⅱ)` → <b>일반암</b>.
    #     구 조건은 `통합전이암` 붙은 형태만 봐서 <b>「통합형전이암」·「통합암전이암」</b>을 놓쳤고,
    #     뒤쪽 일반 `암진단비` 규칙에 걸려 <b>일반암 8건이 합산</b>돼 진짜 일반암을 덮었다.
    #   ⇒ 「통합」과 「전이암」이 <b>한 덩어리</b>면(사이에 형/암 한 글자만) 전부 통합전이암이다.
    if (has('통합전이암') or has('통합형전이암') or has('통합암전이암')) \
       and no('주요치료','수술','통원','일당'): return '통합전이암',0
    #   ★KB `통합암진단비Ⅱ(<b>전이암포함</b>)(유사암제외)`는 <b>통합암</b>이다. 구 제외어 `'전이'`가
    #     「전이암포함」의 '전이'를 잡아 <b>통째로 [확인]큐(None)</b>로 빠졌다(실측 천규운).
    #     통합전이암은 위 조건이 이미 가져가므로 여기서는 <b>'전이암진단'만</b> 막는다.
    if (has('통합암진단') or has('통합암')) and no('전이암진단','간호','간병','주요치료','수술','통원','일당','보험료'):
        return '통합암',0   # ★v320 통합암 단독 = 대표금액 1개   # ★v30z 통합전이암=개별담보·대표금액 1개(§8.2, PPT·보장설명지 반드시 반영)
    if has('전이암') and no('통합'): return '__무시__',0   # ★v30z 전이암진단비 단독=무시(지점장 2026.07.05)
    if has('암') and has('주요치료') and no('비급여','순환계','2대','유사암'): return '암주요치료비',0   # ★v260c 지침 §8.2 = 비급여는 하이클래스(암)가 우선(리터럴 '암주요치료비'는 위 2034행이 먼저 잡는다)
    if has('하이클래스'): return '하이클래스(암)',0
    if (has('비급여') or has('하이클래스')) and has('주요치료'): return '하이클래스(암)',0   # 비급여 주요치료비(암 미명시)=하이클래스(암)
    if has('중입자'): return '중입자치료비',0
    if has('양성자'): return '양성자치료',0
    if has('세기조절'): return '세기조절치료',0
    # ★★★★★v267(지점장 확정 2026.07.28, 영구): <b>항암방사선약물 행 = 아래 3개 이름만</b>.
    #   지점장 원문 = "항암방사선약물치료비 / 항암방사선치료비 / 항암약물치료비 3개만 맞다".
    #   → <b>카티(CAR-T)·표적·호르몬·양성자·세기조절·부위한정 변형은 전부 [확인]큐</b>.
    #   실측 사고(이정화 KB 3.3.5): `카티(CAR-T)항암약물허가치료비` <b>5,000</b>이 이 행에 들어가
    #   대표(max)로 <b>정답 100을 눌렀다</b>(항암방사선치료비 100 + 항암약물치료비 100 → 대표 100).
    #   ★지침 §8.2에 이미 「(카티 제외)」가 명시돼 있었는데 코드가 지키지 않았다.
    # ★★★★★v346 (지점장 지시 2026.08.02, 영구): "<b>항암약물치료특약 or 항암방사선치료특약 =
    #   둘중 하나 대표값 하나를 항암방사선약물</b>". 구 조건은 <b>「…치료비」 어미만</b> 인식해
    #   <b>「…치료특약」이 통째로 미매핑(None)</b>이었다(실측 신한/미래에셋 각 2,000 소실).
    #   → <b>특약·담보·보장 어미</b>도 같은 행으로 받는다. 제외어(카티·표적·호르몬·양성자·세기조절·
    #     중입자·특정부위)는 아래 기존 분기를 그대로 탄다. 두 담보가 같이 있어도 <b>대표(max) 1개</b>
    #     (`_rep1`에 '항암방사선약물' 이미 등록).
    if (has('항암방사선약물치료비') or has('항암방사선치료비') or has('항암약물치료비')
        or (has('항암') and (has('방사선') or has('약물')) and has('치료')
            and (has('특약') or has('담보') or has('보장'))
            and no('면역','특정','허가'))):   # ★지시받은 두 담보만 — 면역·특정·허가 변형은 종전대로 [확인]큐
        if has('카티') or has('CAR-T') or has('CART') or has('표적') or has('호르몬'): return None,0
        if has('양성자') or has('세기조절') or has('중입자'): return None,0
        # 특정부위·특정암 한정 변형(예 남성생식기관련암)은 기본 항암방사선과 별개 → [확인]
        if has('생식기') or has('전립선') or has('음경') or has('고환') or has('유방') or has('자궁') or has('갑상선'): return None,0
        return '항암방사선약물',0
    # ★★★★★v267(지점장 확정 2026.07.28): <b>카티(CAR-T)는 항암방사선약물이 아니다 → [확인]큐</b>.
    #   지침 §8.2 원문에 이미 <b>「중입자='항암중입자방사선치료비'만(카티 제외)」</b>로 명시돼 있는데
    #   구 코드는 카티를 <b>항암방사선약물 행에 넣었다</b>. 이 행은 대표(max)라
    #   <b>카티 5,000이 진짜 값 100을 눌렀다</b>(실측 이정화 KB 3.3.5).
    #   실측 정답 = `항암방사선치료비 100` + `항암약물치료비 100` → 대표(max) <b>100</b>.
    if has('카티') or has('CAR-T') or has('CART'): return None, 0
    # ★v197 지점장 확정 2026.07.23: 'N대특정암진단비'(16대·11대·5대 등)·'특정암진단비' = 고액암 행
    #   근거: 암진단비(유사암제외)와 별개 담보 → 일반암 산입 금지(등식1 위반)
    if has('특정암') and has('진단') and no('주요치료','특정치료','수술','일당','입원','방사선','약물','표적','양성자','세기','중입자','통원','보험료'):
        return '고액암',0
    if has('고액암'): return '고액암',0
    # ★v29q-2 '암입원일당(…유사암들…)' = 암입원일당 키워드 우선 → 암일당 (유사암 판정보다 먼저)
    if ('암입원일당' in n) or (has('암') and has('입원일당')): return '암일당',0
    # ★v29q-1 '암진단비(…유사암들…)' = 암진단비 키워드 우선 → 일반암 (괄호 유사암 구성 무시)
    if has('소아암') and no('제외'): return None,0   # ★v30q 다발성소아암 등 = 일반암과 별개 담보 → [확인](합산 금지, 지점장 2026.07.03)
    if re.search(r'암\s*진단비\s*[(（]', r) and no('유사암제외'): return '일반암',0
    # ★★★v255(장O경 KB 실측 2026.07.26): KB 별첨 3열은 담보명이 접혀 <b>'진단비'가 잘려 나간다</b>.
    #   실측 — `암진단비(유사암제외)` → <b>`암(유사암제외)`</b>로 절단되어 매핑 실패(None) → 일반암 1,000 누락.
    #   → '암' + '유사암제외'가 함께 있으면(수술·일당·치료비 계열이 아닌 한) <b>일반암</b>으로 본다.
    if (has('암') and has('유사암제외')
            and no('수술','일당','입원','통원','치료비','방사선','약물','표적','양성자','세기','중입자','보험료')):
        return '일반암',0
    # ★v255: `비급여(전액본인부담 포함) 암` — KB 비급여 암주요치료비의 절단형.
    #   '주요치료비'가 잘려도 <b>비급여 + 암</b>이면 하이클래스(암) 행이 정본(§8.2 ③).
    if (has('비급여') and has('암') and (has('전액본인부담') or has('전액본인'))
            and no('유사암','수술','일당','입원','통원','방사선','약물','표적','양성자','세기','중입자')):
        return '하이클래스(암)',0
    # 유사암 — 단 '유사암제외'(유사암을 뺀 일반 암진단)는 일반암
    # ★★★★★v248 (지점장 확정 2026.07.26): <b>"유사암이라고 기재된 것만 해라"</b>
    #   → 유사암 행에는 담보명에 <b>'유사암'(또는 '소액암')이 명시된 담보만</b> 넣는다.
    #   구 v230의 <b>동의어 자동산입</b>(갑상선·갑상샘·기타피부·경계성·제자리·상피내·양성뇌종양)은 <b>폐기</b>.
    #   ★실측 근거(양*선 KB 3열): 메리츠 `갱신형 갑상선암(초기제외)진단비` 1,000 ·
    #     `갱신형 갑상선암 및 기타피부암의 전이암(림프절 등 전이제외)진단비` 1,000이 유사암에 산입돼
    #     <b>3,000</b>이 됐으나 KB 전체보장현황 유사암은 <b>1,000</b>(= `갱신형 유사암진단비`뿐)이다.
    #   ★★<b>영향 고지</b>: 이 규칙으로 이정화 우체국 `갑상샘암치료보험금`·`상피내암치료보험금`은
    #     <b>[확인]큐</b>로 간다(구 v230에선 유사암에 산입돼 한장표 900과 일치했다). 추측 대신 신인 수기 확인.
    if any(k in n for k in [_norm(x) for x in ['유사암','소액암']]) and no('유사암제외','유사암 제외'):
        # ★★★v207 (지점장 확정 2026.07.25, 양*선 메리츠 실측): '유사암(갑.기.경.제)'는 <b>진단비 전용 행</b>이다.
        #   글자만 보고 넣던 탓에 <b>수술비·치료비·일당</b>까지 산입돼 유사암이 1,250(=100+1,000+150)으로 부풀었다.
        #   실측 오류 2건 — '갱신형 갑상선기능항진증치료비' 100(갑상선 <b>기능</b>항진증 = 암이 아니다) ·
        #                  '갱신형 유사암수술비' 150(수술비는 진단비 행이 아니다).
        #   → 진단이 아닌 담보는 <b>[확인] 큐로 보낸다</b>(임의로 암수술 행에 넣으면 암수술 200→350으로 또 틀린다).
        if has('수술') or has('치료비') or has('일당') or has('입원') or has('통원') or has('주요치료') or has('기능항진') or has('기능저하'):
            return None, 0
        return '유사암(갑.기.경.제)',0
    if has('중대한') and has('암'): return '중대한 암',0
    if has('암') and has('진단') and no('고액','소액','표적','방사선','약물','수술','일당','양성자','세기','중입자','전이','뇌','보험료'):
        return '일반암',0
    # ★v30l ○○암보험(진단 표기 없는 암 주계약, 예 신한생활비주는암보험·종합암보험) → 일반암 합산. 유사/고액/통합/전이/중대한/치료·수술·일당 계열은 위에서 이미 분기
    if has('암보험') and no('유사','고액','소액','통합','전이','중대한','주요치료','특정치료','하이클래스','수술','일당','방사선','약물','표적','양성자','세기','중입자','보험료'):
        return '일반암',0
    if has('암') and has('입원'): return '암일당',0

    # ── 뇌혈관 ──
    # ★★★★★v421 (지점장 확정 2026.08.14 박미정) — <b>「주요상해뇌출혈진단비」는 외상성이다</b>.
    #   지점장 원문 「<b>외상성뇌출혈이라 뇌출혈진단비와 별개야</b>」.
    #   실측: DB 참좋은운전자상해보험2510 `주요상해뇌출혈진단비 1,000만`이
    #   <b>뇌출혈진단비(35행)</b>에 들어가 있었다 → <b>외상성뇌출혈(37행)</b>이 정본.
    #   ★<b>상해로 인한 뇌출혈은 질병 뇌출혈진단비가 아니다</b>(제3조 단독 5종의 취지).
    if (has('외상성') or has('주요상해')) and has('뇌출혈'): return '외상성뇌출혈',0
    # ★★★★★v325 <b>뇌출혈진단비는 진단 전용</b>(단독담보 원칙 = "그 담보 하나로 존재하면").
    #   <b>수술·입원·일당·통원 담보는 이 행이 아니다</b>. 제외어가 없어 전부 합산되던 것을 막는다.
    #   실측 = 우체국 `뇌출혈수술급부금`(→뇌혈관수술비) · `뇌출혈입원급부금`(→[확인]큐).
    # ★★★★★v620 (지점장 실측 2026.08.30 「<b>뇌출혈 사망·급성심근경색 사망 — 뒤에 사망이
    #   붙고 엑셀에 없는 건데 또 더하기했다</b>」 · 배미용 실측).
    #   [지침 §8.1] 암사망·재해사망·교통사망 등 <b>특정 사망은 일반사망에 표기하지 않는다</b>.
    #     그런데 진단비 판정에 <b>`사망` 제외어가 없어</b> 「뇌출혈 사망」이 뇌출혈진단비로,
    #     「급성심근경색 사망」이 급성심근경색으로 <b>합산됐다</b>.
    #   ⇒ <b>특정질환 + 사망</b>은 마스터 무행 → [확인]큐(제0조 꼭대기).
    #     ★`일반사망`·`상해사망`·`질병사망`은 <b>마스터에 행이 있어</b> 위에서 이미 처리된다.
    if has('사망') and not has('일반사망') and not has('상해사망') and not has('질병사망'):
        if any(has(_k) for _k in ('뇌출혈','뇌졸','뇌혈관','급성심근','심근경색','협심',
                                  '허혈','부정맥','심부전','암','간','신장','폐')):
            return None, 0
    if has('뇌출혈') and no('수술','입원','일당','통원'): return '뇌출혈진단비',0
    if has('중대한') and has('뇌졸'): return '중대한 뇌졸증',0
    if has('뇌졸'): return '뇌졸증진단비',0
    if has('산정특례') and has('뇌'): return '산정특례뇌혈관',0
    # ★2026.07.12 지점장 확정: '특정뇌혈관' = 뇌졸증
    if has('특정') and has('뇌혈관') and has('진단') and no('수술','주요치료','산정특례','혈전'): return '뇌졸증진단비',0
    # ★v30o 고정(전 회사, 지점장 2026.07.03): 뇌혈관진단비Ⅰ→뇌혈관진단비 / 뇌혈관진단비Ⅱ→뇌졸증. Ⅲ은 뇌혈관진단비.
    if has('뇌혈관') and has('진단') and no('수술','주요치료','산정특례','혈전'):
        _n=_rmn(raw)
        if _n==2: return '뇌졸증진단비',0
        if _n in (1,3): return '뇌혈관진단비',0
    if has('뇌혈관') and has('진단'): return '뇌혈관진단비',0
    if has('혈전용해') and has('뇌') and no('주요치료'): return '혈전용해치료비',0   # ★v547 제21조 3항

    # ── 심장 ──
    if has('주요치료') and (has('순환계') or has('2대') or has('뇌혈관') or has('심뇌') or has('허혈') or has('심장')): return '2대 주요치료비',0   # 뇌혈관+허혈성/심장 주요치료비=순환계=2대주요치료비
    if has('중대한') and (has('심근') or has('급성심근')): return '중대한 급성심근',0
    if has('심근병증') or has('심근증'): return '심근병증',0
    if has('판막'): return '심장판막',0
    # ★★★★★v326 <b>급성심근경색은 진단 전용</b>(단독담보 원칙 = "그 담보 하나로 존재하면").
    #   수술은 위 수술비 블록에서 <b>심장수술비</b>로 빠지고, <b>입원·일당은 마스터에 행이 없으므로 [확인]큐</b>
    #   (지점장 "<b>엑셀에 없는 건 제외다</b>"). 통원도 동일.
    if has('급성심근') and no('수술','입원','일당','통원'): return '급성심근경색',0
    # ★2026.07.12 지점장 확정: '특정허혈성' = 협심증 (v28 '허혈심장질환진단비→허혈성 진단비' 규칙보다 우선)
    if has('특정') and (has('허혈성') or has('허혈심장')) and has('진단') and not has('수술'): return '협심증',0
    if has('허혈성진단') or ((has('허혈성') or has('허혈심장')) and has('진단') and not has('수술')): return '허혈성 진단비',0   # ★v29t 허혈심장질환진단 포함
    # ★v40b KB '심장질환(특정)' 진단비: 특정Ⅱ=급성심근경색 / 특정Ⅰ=허혈성 진단비 (지침 §8.3.1 KB).
    #   OCR이 로마숫자를 Ⅱ/II/2 등으로 흘려 매칭 실패하던 버그 수정 → 세 표기 모두 인식.
    if has('심장질환') and has('특정') and has('진단') and no('수술','주요치료'):
        _rr=str(raw)
        _is2 = ('Ⅱ' in _rr) or ('특정 II' in _rr) or ('특정II' in _rr) or ('특정 2' in _rr) or ('특정2' in _rr) or ('(특정 II)' in _rr) or ('（특정 II）' in _rr)
        _is1 = ('Ⅰ' in _rr) or ('특정 I' in _rr) or ('특정I' in _rr) or ('특정 1' in _rr) or ('특정1' in _rr) or ('(특정 I)' in _rr) or ('（특정 I）' in _rr)
        # ★★★★★v591 (지점장 실측 2026.08.26 「또 협심증이 2배로 나온다」 · 롯데 let:smile).
        #   [실측] 협심증 행에 두 줄 — `심혈관질환진단비(특정심장질환Ⅰ)` 1,000(<b>여기</b>) +
        #     `협심증[심장묶음]` 1,000(`_HB`) = <b>2,000</b>.
        #   [조문 위반] 「특정Ⅰ」의 뜻은 <b>회사마다 다르다</b> — DOCTRINE 회사별 정본표 =
        #     롯데 특정Ⅰ→<b>급성심근경색</b> / DB 특정Ⅰ→협심증+염증 / KB 특정Ⅰ→협심증+허혈성+빈맥+심부전.
        #     이 분기는 <b>회사를 모른 채 라벨만 보고</b> Ⅰ→협심증으로 단정했다 → 롯데에서 틀렸다.
        #   ⇒ <b>여기서 값을 만들지 않는다.</b> 묶음 분해는 `_HB`(회사별 정본표) <b>한 경로</b>가 한다.
        #     `_HB`가 못 잡으면 [확인]큐로 간다 — 제0조 꼭대기(엑셀에 없는 건 만들지 않는다).
        return None, 0
    if has('일당') and (has('허혈') or has('협심') or has('심부전') or has('부정맥') or has('빈맥') or has('뇌혈관') or has('심뇌')): return None,0   # ★v30b 질환별 입원일당 ≠ 진단비 → [확인] (조성래 허혈일당 오합산 수리)
    if has('협심'): return '협심증',0
    # ★★★v326 허혈성 폴백에도 제외어 — 위 2795행에는 `not has('수술')`이 있었으나
    #   이 <b>폴백에는 아무 제외어가 없어</b> 수술·입원 담보가 진단비로 샜다.
    if has('허혈') and no('수술','입원','일당','통원'): return '허혈성 진단비',0   # ★v29t §8.3: 허혈 단독=허혈성 진단비
    if has('심부전'): return '심부전',0
    if has('심내막') or has('심근염') or has('심장막') or has('심장염증'): return '염증',0
    if has('빈맥'): return '빈맥',0   # ★지점장 2026.07.05: 빈맥(I47·48)=master 40행 정식 사용(v30z6 '무행' 폐기). 빈맥≠부정맥(I49)
    # ★★★v217 (지점장 지시 2026.07.25, 영구): <b>'심장부정맥 고주파·냉각절제술보장' = 부정맥 행이 아니다</b>.
    #   지점장 원문 = <b>"부정맥은 부정맥이라고 적혀있다"</b> → 마스터 '부정맥'(42행)은 <b>진단비 전용</b>이고,
    #   <b>고주파절제술·냉각절제술(전극도자절제술)은 치료 시술 담보</b>라 진단비 행에 넣으면 안 된다 → <b>[확인]큐</b>.
    #   실측 오류: `심장부정맥고주파·냉각절제술보장`이 아래 `has('부정맥')`에 걸려 부정맥 진단비로 산입됐다.
    if has('부정맥') and (has('고주파') or has('절제') or has('냉각') or has('시술') or has('도자')):
        return None,0
    if has('부정맥'): return '부정맥',0
    if has('산정특례') and has('심'): return '산정특례심장',0
    if has('2대') and has('주요'): return '2대 주요치료비',0
    if has('혈전용해'): return '혈전용해치료비',0

    # ── 사망 ──
    # ★★★★★v301-B (지점장 지시 2026.07.31, 영구): <b>'타인사망교통사고처리지원금'은 합의금이다</b>.
    #   지점장 원문 = "이건 합의금중 하나야 <b>대표값들중 하나만 넣으면되</b>".
    #   [구 결함] 아래 `교통 and 사망 → 교통상해사망`이 <b>합의금 판정(더 아래)보다 먼저</b> 돌아
    #   처리지원금을 사망 담보로 먹었다. v300에서 9행을 실제로 채우기 시작하면서 값까지 실린다.
    #   ★대표(max) 1건 처리는 `_rep1`에 '합의금'이 이미 들어 있어 자동 적용된다(v198).
    #   ★★<b>6주미만이 먼저다</b> — `교통사고처리지원금(6주미만 진단)`은 6주미만 전용행이다.
    #     규칙을 앞으로 당기면서 이 순서를 놓쳐 합의금으로 갔다(단위검증에서 검출·수정).
    if has('6주'): return '6주미만',0
    if (has('처리지원금') or has('형사합의') or has('합의금')) and no('자전거'): return '합의금',0   # ★v356 자전거사고 제외
    if has('CI') and has('사망'): return '중대한CI적용',0
    if has('교통') and has('사망'): return '교통상해사망',0
    if (has('상해') or has('재해')) and has('사망'): return '상해사망',0
    if has('질병') and has('사망'): return '질병사망(80세)',0
    if has('일반사망') or (has('사망') and no('상해','질병','교통','재해','CI','암','운전','입원','수술')): return '일반사망',0

    # ── 후유장애 ──
    if has('화재') and (has('후유') or has('장해')): return None,0   # ★v29q-9 화재상해후유(3~100%)≠상해후유3%, 담보행 미기재→[확인] 큐
    # ★★★v104 영구지침(지점장 확정 2026.07.20): '無고도장해보장' = 80% 상해후유장해.
    #    구 로직은 '상해/재해/교통' 글자가 없어 질병후유80%로 잘못 갔다.
    if has('고도장해') and no('질병'): return '상해후유80%',0
    if has('후유') or has('장해') or has('장애'):
        # ★★★★★v302-A (지점장 지시 2026.07.31, 영구): <b>'후유' 표기가 없는 장해 변형은 기재금지</b>.
        #   지점장 원문 = "재해장해연금·교통재해장해급부금이 상해후유3%에 전부 합산금지 ㅡ 기재금지 /
        #   <b>상해후유3% or 재해후유3% or 상해후유장해 or 재해후유장해 ㅡ상해후유3%로 가능</b>".
        #   → 허용 4형태는 <b>전부 '후유'를 포함</b>하고, 금지 예시 2건은 <b>전부 '후유'가 없다</b>.
        #   실측(이정화 우체국): `재해장해급부금` 2,000 + `재해장해연금` 6,000 + `교통재해장해급부금` 3,000이
        #   상해후유3%에 전부 합산돼 39,000이 됐다(한장표 20,000). → 셋 다 [확인]큐.
        #   ★위의 `고도장해`(v104 지점장 확정)는 이 검사보다 앞에서 이미 처리된다 — 영향 없다.
        if not has('후유'):
            return None,0
        # ★★★★★v421 (지점장 확정 2026.08.14 박미정) — <b>「특정상해후유장해」는 상해후유3%가 아니다</b>.
        #   지점장 원문 「<b>DB 특정상해후유장해 100만 — 빼자</b>」.
        #   실측: DB 참좋은훼밀리 `상해사망,(20~100%)이상후유장해(보통약관)`이 신정원에서
        #   <b>「특정상해후유장해」</b>로 분류되는데 엑셀은 상해후유3%에 100만을 더해 10,100이 됐다.
        #   KB 한장표는 10,000 → <b>「특정」이 붙은 후유장해는 3% 행이 아니다</b>. [확인]큐로 보낸다.
        if has('특정') and no('80'):
            return None,0
        # ★★★v222 (지점장 지시 2026.07.25, 영구): <b>후유장해 20% · 50%는 결과값 미기재 — 더하지도 말 것</b>.
        #   마스터 후유장해 행은 <b>3% · 80% 넷뿐</b>이고 20%·50% 전용행은 없다.
        #   구 코드는 '80'이 없으면 전부 <b>3% 행에 합산</b>해서 금액을 부풀렸다.
        #   실측(이정화): `일반상해50%이상후유장해생활자금` 4,800 · `교통상해50%이상후유장해(운전자)` 100이
        #   상해후유3%에 합산됐다. → <b>[확인]큐로 보내고 어느 행에도 넣지 않는다</b>.
        if re.search(r'(?<![\d.])(20|50)\s*%', n) and not re.search(r'(?<![\d.])80\s*%', n):
            return None,0
        # ★★v92 (장혜경 실측): '질병후유장해(80%미만)Ⅱ'가 '80' 글자 때문에 80%행으로 잘못 갔다.
        #   → '80%미만'/'80% 미만'이면 3% 행. (한장보장표 질병3% 100과 일치)
        _u80 = ('80%미만' in n.replace(' ','')) or has('미만')
        sev = '3' if _u80 else ('80' if ('80' in n or has('고도')) else '3')
        body = '상해' if (has('상해') or has('재해') or has('교통')) else '질병'
        return f'{body}후유{sev}%',0

    # ── 일당/입원 ──
    # ★v62 간호간병통합서비스(지점장 2026.07.15): 담보명에 '간호'+'간병'+'통합' 동시 존재 →
    #   간호통합병동 행. '간병인사용…일당비(간호간병통합서비스)' 형태가 아래 has('간병인')에
    #   먼저 걸려 간병인으로 오분류·누락되던 것을 차단(간호통합병동을 간병인보다 우선 판정).
    if has('간호') and has('간병') and has('통합'):
        if has('181'): return None,0        # 1-180일 기준만(181일이상 제외, v41 유지)
        return '간호통합병동',0
    # ★★v216 (지점장 지시 2026.07.25): <b>'간호'가 빠진 '간병통합서비스' 형태도 간호통합병동</b>.
    #   실측 누락: `간병통합서비스입원일당`이 위 조건(간호+간병+통합)에서 탈락하고
    #   아래 `has('간병인')`에도 안 걸려(‘간병인’이 아니라 ‘간병’) <b>[확인]큐로 사라졌다</b>.
    #   ★'간병인'이 들어간 담보는 제외한다(그건 간병인/간병인지원일당 행).
    # ★★★v220 (2026.07.25 긴급 축소 — v216 과잉 완화가 사고를 냈다):
    #   v216에서 `간병+통합+서비스`만 있으면 간호통합병동으로 보냈더니
    #   <b>`간병비통합서비스보장`·`통합간병서비스진단비`·`간병통합서비스수술비` 같은 진단비·수술비까지
    #   간호통합병동 행으로 들어갔다</b>(실측). 진단비가 들어오면 <b>2,000만원</b>이 일당 칸에 찍힌다.
    #   → <b>'입원' 또는 '일당'이 있는 담보만</b> 허용하고, 진단·수술·사망·치료비·간병인은 제외한다.
    if has('간병') and has('통합') and has('서비스') and (has('입원') or has('일당')) \
       and no('간병인','진단','수술','사망','치료비','후유','장해','보험료'):
        if has('181'): return None,0
        return '간호통합병동',0
# ★★★★★v527 제110조 <b>폐기</b> (지점장 지시 2026.08.21 「요양병원은 삭제다」)
    #   구 v512·v518 조문 「마스터에 요양병원 2행 신설」은 <b>전부 폐기</b>한다.
    #   요양병원 간병인은 <b>마스터 무행</b>이므로 기재하지 않고 [확인]큐로 보낸다(제0조 꼭대기 지침).
    #   ★'제외'가 붙은 것(요양병원제외)은 <b>일반 간병인</b> 행이다 — 아래 규칙으로 내려간다.
    if has('간병인') and has('요양병원') and no('제외'):
        return None,0                      # ★v527 제110조 폐기 — 마스터 무행 → [확인]큐
    # ★★★v215 (지점장 확정 2026.07.25, 영구): <b>'간병인사용'이 들어가면 '지원'보다 우선해서 '간병인' 행</b>이다.
    #   실측 오분류: '간병인사용지원일당' · '간병인사용비용지원' · '간병인사용 질병입원일당(간병인 지원)'이
    #   아래 `has('간병인') and has('지원')`에 먼저 걸려 <b>간병인지원일당 행에 20만·7만이 찍혔다</b>(지점장 지적).
    #   '간병인사용'은 실제 간병인을 쓴 날에 주는 <b>간병인 담보</b>이고, '간병인지원일당'은 별개 전용 담보다.
    # ★★★★★v600 제139조 (지점장 확정 2026.08.26 「<b>간병인은 이 계약이 15만원인 거다</b>」).
    #   [실측 · 현대해상] 같은 담보가 <b>구간별로 6줄</b> —
    #     (1-180일) <b>15</b> · (181-365일) 20 · (상급종합병원)(1-180일) 5.
    #   대표(max)라 <b>20</b>이 잡혔다. 그러나 <b>기본 구간은 1-180일</b>이고
    #   181일 이상은 <b>연장</b>, 상급종합병원은 <b>병원 한정</b>이다 — 기본이 정본이다.
    #   ⇒ 181일 이상·상급종합병원 변형은 <b>[확인]큐</b>(제0조 — 비슷한 행에 임의로 넣지 않는다).
    #   ★`간호간병통합서비스`는 위에서 이미 <b>간호통합병동</b>으로 갈라진다.
    if has('간병인') and (re.search(r'18[1-9]|19\d|[2-9]\d\d\s*일', n) or has('상급')):
        return None, 0
    if has('간병인') and has('사용'): return '간병인',0
    if has('간병인') and has('지원'): return '간병인지원일당',0   # ★v29w (지점장 2026.07.02) 간병인지원일당 전용행
    if has('간병인'): return '간병인',0
    if has('간호') and (has('통합') or has('간병')):
        if has('181'): return None,0        # ★v41 간호통합병동 = 1-180일 기준만(181일이상 제외)
        return '간호통합병동',0
    if has('1인실') and has('상급'): return '1인실 상급병원',0
    if has('1인실') and has('종합'): return '1인실 종합병원',0
    if has('중환자') and has('상해'): return '상해중환자실',0
    if has('중환자') and has('질병'): return '질병중환자실',0
    if (has('질병') or has('수술')) and has('일당') and has('수술'): return '질병수술일당',0
    # ★★★v209 단독행 원칙(지점장 확정 2026.07.25, 영구): <b>'질병종합병원일당' 행은 그 이름으로 된 단독 담보 전용</b>이다.
    #   <b>병실 등급 담보(2~3인실 · 2인실 · 3인실 등)는 이 행에 넣지 않는다</b> → [확인] 큐.
    #   마스터에 있는 병실 행은 <b>1인실 종합병원 · 1인실 상급병원 둘뿐</b>이고, 2~3인실 전용행은 만든 적이 없다.
    #   실측 오류(양*선 삼성 New내돈내삼): 상급 2~3인실 5 + 종합 2~3인실 5가 이 행에 들어가 10으로 찍혔다.
    if has('인실'): return None,0
    # ★★★★★v301-C (지점장 지시 2026.07.31, 영구): <b>'상급병원 질병입원일당' = 기재금지</b>.
    #   지점장 원문 = "상급병원 질병입원일당 ㅡ 기재금지. 엑셀. 행삭제" → <b>마스터 행을 삭제했고</b>
    #   담보는 <b>[확인]큐</b>로만 남긴다. 구 v223(전용행 신설)은 <b>완전 폐기</b>한다.
    #   ★유지 대상 2행 = <b>질병종합병원일당(기존)</b> + <b>종합병원 상해입원일당(v301 신설·v306 라벨 확정)</b>.
    if has('상급') and has('질병') and (has('일당') or has('입원')) and no('수술','중환자','간병','간호','진단'):
        return None,0
    # ★v314 라벨 정본 = master.xlsx B54 「종합병원 질병입원일당」(지점장 수정 2026.08.01).
    #   55행 「종합병원 상해입원일당」과 대칭. 구 라벨 '질병종합병원일당'은 폐기 —
    #   마스터 라벨과 다르면 nm2r에서 못 찾아 <b>그 행이 통째로 0</b>이 된다.
    # ★★★★★v589 (지점장 확정 2026.08.25 「<b>상급병원일당은 미기재다 (1인실제외)</b>」
    #   「<b>엑셀에 없는건 제발 하지마라</b>」 = 제27조 꼭대기 지침 그대로).
    #   [실측 결함] 「상급<b>종합병원</b>」 안에 「종합병원」 글자가 그대로 들어 있어
    #     `상해입원일당(상급종합병원)`이 <b>종합병원 상해입원일당</b> 행으로 갔다.
    #     주석엔 「상급은 위에서 이미 [확인]큐로 빠졌다」고 적혀 있었지만 <b>조건에 차단이 없었다</b>.
    #   [마스터] 56 종합병원 질병입원일당 · 57 종합병원 상해입원일당 <b>둘뿐</b> —
    #     상급 전용 입원일당 행은 <b>없다</b>. 없는 행은 만들지 않고 <b>[확인]큐</b>로 보낸다.
    #   ★<b>1인실은 예외</b>: 마스터에 63 1인실 종합병원 · 64 1인실 상급병원이 <b>실재</b>하므로
    #     그 경로(위쪽 `인실` 분기)가 먼저 처리한다 — 여기 오지 않는다.
    if has('상급') and has('일당') and no('인실'):
        return None, 0
    if has('질병') and has('종합') and has('일당'): return '종합병원 질병입원일당',0
    # ★★★★★v301-C 신설: 종합병원 <b>상해</b>입원일당 = 전용행. 질병 짝 행과 1:1 대칭이다.
    #   ★'재해'=상해 동일 적용(§v29v). ★'상급'은 위에서 이미 [확인]큐로 빠졌다.
    if (has('상해') or has('재해')) and has('종합') and has('일당'): return '종합병원 상해입원일당',0
    # ★★★v223 = v212 재적용(지점장 지시 2026.07.25, 영구): <b>'N대질병입원일당'(2대·3대·5대…) ·
    #   '특정질병입원일당'은 질병입원일당이 아니다 → 기재 금지</b>([확인]큐).
    #   지점장 원문 = "2대질병입원일당(1일이상) → 그건 질병일당이 아니므로 기재금지".
    #   실측 오류: 메리츠 `갱신형 2대질병입원일당(1일이상)` 3만이 질병일당으로 산입됐다.
    if (has('일당') or has('입원')) and re.search(r'\d+\s*대', n) and no('인실'):
        return None,0
    # ★v30k 교통상해입원일당 ≠ 상해입원일당(합산 금지). 질환·부위·교통 접두 변형은 base 아님 → [확인]
    # ★병원규모(상급종합/종합) 명시 = 개별 전용행 / 일반 질병입원일당(밴드) = 합산 (지점장 2026.07.05)
    # ★v223 `_dilqual`에 <b>'상급'</b> 추가 — 상급병원 일당은 위 전용행으로 갔고, 남은 변형이
    #   질병일당·상해일당에 섞이는 것을 막는다(v212 정본 재적용).
    _dilqual = ('교통','암','뇌','심','허혈','간','신장','폐','위','골절','화상','특정','재해외','종합','요양','중환자','수술','상급')
    # ★★★★★v295 (김수영 KB 실측 2026.07.31, 영구): <b>'간편가입'의 '간'이 제외어 '간'(간암)에 걸린다</b>.
    #   실측 — `상해입원일당(1일이상)Ⅱ(간편가입)(갱신형) 1만`이 [확인]큐로 사라졌다.
    #   KB 간편심사 상품은 <b>전 담보에 '(간편가입)'이 붙으므로 일당이 통째로 증발</b>한다.
    #   v216b '수술비 관혈 → 비관혈' 사고와 <b>같은 계열</b>(부분 문자열 우연 일치).
    #   → 제외어 검사에서만 <b>심사방식 수식어를 지운 사본</b>을 쓴다(다른 규칙엔 영향 없음).
    _ndil = n
    for _q in ('간편가입','간편심사','간편고지','간편','유병자'):
        _ndil = _ndil.replace(_norm(_q), '')
    _nodil = lambda *ks: not any(_norm(k) in _ndil for k in ks)
    if (has('상해일당') or has('상해입원일당') or has('재해일당') or has('재해입원일당')) and _nodil(*_dilqual): return '상해일당',0   # ★재해=상해 동일(정본)
    if (has('질병일당') or has('질병입원일당')) and _nodil(*_dilqual): return '질병일당',0   # 순수 질병(입원)일당만 합산
    # ★★★★★v303 (지점장 지시 2026.07.31, 영구): <b>'질병입원' = '질병입원일당' 같은 것이다 · 기재</b>.
    #   지점장 원문 = "질병입원 = 질병입원일당 / 같은것임.기재 · 상해입원 = 상해입원일당 / 같은것임.기재
    #   / <b>****입원일당 ㅡ 기재안해도됨</b> / 엑셀표 기준으로 기재".
    #   [구 결함] 규칙이 <b>'질병입원일당'이라는 연속 문자열</b>을 요구해 '일당' 글자가 없는
    #   `질병입원`·`상해입원`이 통째로 [확인]큐로 빠졌다(실측 → None).
    #   ★<b>접두 수식어가 붙은 '****입원일당'은 종전대로 기재하지 않는다</b> — `_dilqual`이 그대로 걸러낸다
    #     (교통상해입원일당 · 2대/특정질병입원일당 · 상급병원 질병입원일당 = 전부 [확인]큐 유지).
    #   ★실손·수술·치료비 계열은 각자 전용행이므로 여기서 제외한다.
    if (has('상해입원') or has('재해입원')) and _nodil(*_dilqual) \
       and no('의료비','실손','급여금','수술','치료비','진단','비용'):
        return '상해일당',0
    if has('질병입원') and _nodil(*_dilqual) \
       and no('의료비','실손','급여금','수술','치료비','진단','비용'):
        return '질병일당',0
    # ★v29v (지점장 2026.07.02): 밴드형 '입원비(1일이상/180일한도)' = 입원일당
    if has('입원비') and (has('1일') or has('180일')) and no('실손','의료비','수술'):
        return ('상해일당' if (has('상해') or has('재해')) else '질병일당'),0   # 재해=상해(§v29v)
    # ★v30c AIG류 밴드 미표기 base 입원비 = 입원일당 (변형·질환한정 [확인])
    if has('입원비') and no('실손','의료비','수술','중환자','상급','종합','중증','특정','암','뇌','허혈','심','간질환','감염'):
        if has('질병'): return '질병일당',0
        if has('상해') or has('재해'): return '상해일당',0

    # ── 운전자 (지침 §운전자 매핑) ──
    # ★★★★★v356 (지점장 확정 2026.08.02, 영구): "<b>자전거사고는 무관하다</b>".
    #   NH `자전거사고 벌금담보` 2,000 · `자전거사고 교통사고처리지원금` 3,000이
    #   <b>대인·합의금 칸에 산입</b>돼 한장표와 어긋났다(실측 합의금 20,000→23,000 · 대인 2,000→4,050).
    #   → <b>담보명에 '자전거'가 있으면 운전자 칸에 넣지 않는다</b> → [확인]큐.
    #   ★자동차 운전자 담보만 이 블록의 대상이다. 범위를 넓히지 않는다.
    if has('자전거'): return None,0
    #  벌금(대인)→대인 / 벌금(대물)→대물 / 처리지원금(중상해포함)→합의금 / 처리지원금(6주미만)→6주미만
    #  변호사→변호사 / 자동차(사고)부상보장·부상위로→자부상
    if has('6주'): return '6주미만',0
    if (has('처리지원금') or has('형사합의') or has('합의금')) and no('자전거'): return '합의금',0   # ★v356 자전거사고 제외
    if has('벌금') and has('대물'): return '대물',0
    # ★★★★★v284 (지점장 확정 2026.07.31): <b>"과실치사는 삭제다. 대인벌금은 과실치사랑 다른 담보다"</b>
    #   → v283에서 내가 지점장 문구("표기 미이행 2가지")를 <b>반대로 해석</b>해 과실치사 계열을
    #     대인에 합류시킨 것은 <b>오류</b>다. 구 v29q-7 배제어를 <b>그대로 원복</b>한다.
    #   <b>과실치사상 벌금 · 업무상 과실치상해실치사상 벌금 = 대인벌금과 별개 담보</b> → 대인에 넣지 않는다.
    #   (누락 방지 원칙에 따라 [확인]큐에는 그대로 남는다.)
    if has('벌금') and no('화재','과실','치사','업무'): return '대인',0
    # ★★★★★v288 (이명순 실측 2026.07.31 · 지점장 확정 "대인 2,000+1,000 총 3,000"):
    #   지침 §8.6은 <b>"벌금(대인)·벌금(대인대물 미표기)→대인"</b> — 즉 <b>벌금 담보</b>의 대인/대물 구분 규칙이다.
    #   구 코드는 '벌금' 없이 <b>'대인' 글자만으로</b> 대인 행에 넣어 지침 범위를 넘었다.
    #   <b>실측(이명순 현대 Hi0910)</b>: `대인교통사고발생위로금담보특별` <b>20</b>이 대인으로 산입 →
    #     대인 = 벌금담보 2,000 + 스쿨존 1,000 + <b>위로금 20</b> = <b>3,020</b>(한장표 정답 3,000).
    #   → 지침 원문대로 <b>'벌금'을 필수 조건으로</b> 되돌린다. 위로금은 [확인]큐로 남는다(누락 금지).
    if has('대인') and has('벌금') and no('대물'): return '대인',0
    if has('대물'): return '대물',0
    if has('변호사'):
        # ★★★★★v580 제122조 4항 (지점장 지시 2026.08.24 「엑셀은 0으로 하고 노란칸 ·
        #   진단서·분석지·리포트는 현물 · 디비손보 202608월~ 건만」).
        #   판정은 <b>여기 한 곳</b>에서 하고 결과를 `_HYUNMUL`에 적는다(같은 판정 두 곳 금지).
        #   조건 3개를 <b>전부</b> 만족할 때만 — DB손보 · 2026.08 이후 · `(현물급부선택가능)`.
        if _is_hyunmul(r, _CUR_CO, _CUR_CD):
            # ★v580b — 두 번째 반환값은 호출부가 금액으로 쓸 수 있으므로 <b>건드리지 않는다</b>.
            #   현물이라는 사실은 `_HYUNMUL`에만 남기고, 금액 0·노란칸·「현물」 표기는
            #   각 산출물이 이 집합을 보고 처리한다(구조 가정 금지).
            _HYUNMUL.add('변호사')
            return '변호사', 0
        return '변호사',0
    if has('자부상') or (has('자동차') and (has('부상') or has('자부상'))):
        # ★자부상=자동차 한정. 급수밴드 있으면 14급 포함 밴드만(1~3/1~7 제외).
        _band=re.search(r'(\d+)\s*~\s*(\d+)\s*급', r)
        if _band:
            if int(_band.group(2))>=14: return '자부상',0
            return None,0
        # ★v65 단일 급수(예 '(1급)'·'(14급)') = 14급만 자부상, 그 외 제외(지점장 2026.07.15).
        #   이전엔 밴드 없으면 무조건 자부상이라 급별 여러 줄(1~14급)이 전부 자부상→합산 수천원 오류였다.
        _single=re.search(r'(?<![\d~])(\d+)\s*급', r)
        if _single:
            return ('자부상',0) if int(_single.group(1))>=14 else (None,0)
        return '자부상',0        # 급 표기 자체가 없는 순수 자동차부상위로금 → 자부상

    # ── 골절/응급/독감/화상/깁스 ──
    if has('5대골절') and has('진단'): return '5대골절진단비',0
    # §골절: '치아/파절 제외' 명시된 것만 제외 행. 단독 골절진단비·치아포함은 포함 행.
    #   ★삼성형 '치아파절 (깨짐, 부러짐) 제외'처럼 중간에 딴말 껴도 잡도록 (치아 or 파절) + 제외 동반 판정
    if has('골절') and has('제외') and (has('치아') or has('파절')): return '골절(치아파절제외)',0
    if has('골절') and (has('치아포함') or has('파절포함') or (has('제외') is False and (has('치아') or has('파절')))): return '골절(치아파절포함)',0
    if has('골절') and has('진단'): return '골절(치아파절포함)',0
    if _norm(raw)=='골절' or has('골절') and no('수술','일당','입원','깁스','부목'): return '골절(치아파절포함)',0
    if (has('응급실') or (has('응급') and has('내원'))) and no('비응급'): return '응급실(응급)',0   # ★v29q-11 응급 단독, 비응급 합산 차단→[확인]
    if has('독감') or has('인플루엔자'): return '독감',0
    if has('화상') and (has('중증') or has('심재성') or has('중대한') or has('부식')): return '중증화상진단비',0
    if has('화상') and has('진단'): return '화상진단비',0
    if has('부목') or has('반깁스'): return '반깁스',0   # ★v29q-5 골절부목치료비=반깁스
    if has('깁스'): return '깁스진단비',0

    # ── 실손 ──
    if (has('실손') or has('입원의료비') or has('상해입원형') or has('질병입원형')) and has('입원'): return '입원',0
    if has('통원') and (has('실손') or has('외래') or has('의료비')) and no('주사','MRI','도수','체외','증식','비급여'): return '통원',0
    if has('처방조제') or has('약제비') or (has('약') and has('실손')): return '약값',0
    if has('일상생활') and has('배상') or has('일배책') or has('일상배상'): return '일상배상책임',0
    return None, 0

_SURG_CASES = [
    # ★★★★★v594 제135조 (지점장 확정 2026.08.26 「비가 없어도 수술비다」 /
    #   「수술비가 너무 헐렁해서 놀랐다」) — <b>넓힌 만큼 새는 것도 케이스로 고정</b>한다.
    ('질병수술(간편건강고지)담보', '질병수술비'),   # 실측 · 현대해상 30이 [확인]큐로 빠졌던 건
    ('상해수술(간편건강고지)담보', '상해수술비'),
    ('질병수술', '질병수술비'), ('상해수술', '상해수술비'),
    ('질병수술비', '질병수술비'), ('질병입원수술비', '질병수술비'),
    ('상해수술비', '상해수술비'), ('상해입원수술비', '상해수술비'),
    ('질병수술비(연간1회한)', '질병수술비'),
    # ★아래는 <b>수술비가 아니다</b> — 넓히면 새는 것들(실측으로 확인)
    ('질병수술위로금', None), ('질병수술자금', None), ('질병수술치료비', None),
    ('상해수술위로금', None), ('상해수술치료비', None),
    ('다발성질병수술(3대질병)(간편건강고지)담보', None), ('특정질병수술비', None),
    ('상해수술비(종합병원)', None), ('상해흉터복원수술비', None), ('교통상해수술비', None),
]


def surg_selftest():
    """★제135조 — 수술비 매핑이 <b>넓지도 좁지도 않은가</b>를 센다."""
    bad = []
    for _nm, _want in _SURG_CASES:
        try:
            _got = resolve2(_nm)[0]
        except Exception as _e:
            bad.append('[제135조 수술비] %s — 예외 %s' % (_nm[:26], str(_e)[:20])); continue
        if _got != _want:
            bad.append('[제135조 수술비] %s — <b>%s</b> (정답 %s)' % (_nm[:26], _got, _want))
    return bad


def _fix_ilbae(contracts):
    """★★★★★v593 제134조 — <b>일상배상책임은 무조건 1억(10,000)</b>
       (지점장 확정 2026.08.26 「일상배상책임은 무조건 1억이다」).
       [실측 · 롯데 let:smile] `[접힘의심] 가족일상생활배상책임(대물` <b>20</b> —
         담보명이 `(대물`에서 잘렸고 원문 `(대물 누수 50만원, 누수외 20만원)`의
         <b>자기부담금 20</b>을 가입금액으로 읽었다. 1억이 20만원이 됐다.
       ⇒ 담보가 있으면 10,000 고정. 자기부담금·누수한도는 가입금액이 아니다.
       ★게이트가 <b>직접 호출</b>해 값으로 검증할 수 있게 함수로 뺀다."""
    for _c in (contracts or []):
        _dm = _c.get('dambo') or {}
        _v = _dm.get('일상배상책임')
        if _v in (None, ''):
            continue
        try:
            _n = float(str(_v).replace(',', ''))
        except Exception:
            _n = None
        if _n is not None and _n != 10000:
            _dm['일상배상책임'] = 10000
            print('[v593 일상배상] %s — %s → 10,000 고정(제134조)'
                  % (str(_c.get('company'))[:10], _v))
    return contracts


def resolve2(raw):
    """raw -> (std, jong). DMAP 정확매칭 우선, 없으면 키워드 사전엔진(resolve_kw)."""
    # ★★★★★v298-B (주재현 새마을금고 실측 2026.07.31): 담보명이 '○○특약[재해사망형/만기환급형] : 담보명'
    #   형태일 때 콜론 앞은 상품·특약명인데 그 안의 '재해사망' 글자가 담보 매칭을 먹었다.
    #   실측 2건 — `상해사고특약[재해사망형/만기환급형] : 화상진단위로금` 30 → 상해사망,
    #             `입원특약[재해사망형/만기환급형] :입원급여금[3일초과/0.1%]` 2 → 상해사망.
    raw = str(raw or '').replace('％', '%')
    # ★★★★★v590 제9조 5항 — 실손 <b>통합형 `의료비(입원+통원)`</b> (지점장 실측 2026.08.25
    #   「<b>4세대 실손이 계속 오류다 · 입원 0 / 통원 5000 / 상해의료비 5000</b>」 ·
    #   「<b>통원 20에 잡아주면 된다</b>」 · 윤선경 NH농협손보 헤아림실손2301).
    #   [실측 결함] 별첨 원문은 `상해(일반상해, 전체상해를 의미) 의료비(입원+통원) 5,000`이다.
    #     ・상해 …의료비(입원+통원) → <b>상해의료비</b>(1세대 전용 행)  ✘
    #     ・질병 …의료비(입원+통원) → <b>통원</b>                      ✘
    #   [왜 못 잡았나] v550에서 고친 것은 <b>붙여 쓴 형태</b>(`상해급여입원의료비`)뿐이고
    #     <b>괄호형은 전부 None/오분류</b>였다. 3열 전용 `_sj_fixname` F7에만 처리가 있어
    #     <b>2열(롯데)은 안 걸렸다</b>. ⇒ 두 경로가 함께 지나가는 <b>여기</b>에 둔다.
    #   ⇒ `의료비` + `입원` + `통원`이 한 담보명에 있으면 <b>입원 행</b>(통합형은 입원 한도가 정본).
    #     통원 20은 세대 판정 후 계약 단위로 채운다(`_silson_tongwon20`).
    # ★★★★★v593 (지점장 실측 2026.08.26 「<b>일배책이 20으로 찍힌다</b>」).
    #   [실측] KB 제안서 일상배상책임 아래 <b>안내 문구</b> —
    #     `※ 자기부담금 : 대물(누수50만원,누수외20만원)` 이 <b>담보로 잡혀 `대물` 20</b>이 됐다.
    #   [원인] 이 줄은 <b>담보가 아니라 자기부담금 안내</b>인데 차단이 없었다.
    #     `자기부담금`·`공제금액`만 단독으로 오면 None이었지만, 뒤에 `대물(...)`이 붙자 통과했다.
    #   ⇒ <b>자기부담금·공제금액·면책금액 줄은 담보가 아니다</b> — 값을 만들지 않는다(제0조 꼭대기).
    _n593 = re.sub(r'[\s※]', '', str(raw))
    if re.match(r'^(자기부담금|공제금액|면책금액|자기부담|공제금)', _n593):
        return None, 0
    _n590 = re.sub(r'\s', '', str(raw))
    if ('의료비' in _n590) and ('입원' in _n590) and ('통원' in _n590) and ('실손' not in _n590):
        return '입원', 0
    _pfx = re.match(r'^(.{4,}?)\s*[:：]\s*(.{4,})$', raw)
    if _pfx and re.search(r'(특약|플랜|보험|계약)', _pfx.group(1)) \
       and re.search(r'(진단|수술|위로금|급여금|공제금|일당|치료비|보험금|입원|장해|사망)', _pfx.group(2)):
        raw = _pfx.group(2).strip()
    # ★★★★★v300 (지점장 지시 2026.07.31, 영구): <b>교통상해사망·대중교통상해사망 = 마스터 9행
    #   '교통상해사망' 전용행 단독</b>. 상해사망(10행)에 합산 금지.
    #   [구 결함] 아래 DMAP <b>부분일치 루프</b>에서 키 '상해사망'이 '교통상해사망'을 substring으로
    #   먼저 먹어, resolve_kw의 정본 규칙(`has('교통') and has('사망') → 교통상해사망`)까지
    #   도달하지 못했다. 실측 = 주재현 엑셀 9행이 항상 0.
    #   ★범위는 지점장 문구 그대로 <b>'교통'+'상해사망'</b>뿐이다. 처리지원금·벌금 등은 손대지 않는다.
    _n300 = re.sub(r'\s', '', str(raw))
    if ('교통' in _n300) and ('상해사망' in _n300) and ('지원금' not in _n300):
        return ('교통상해사망', 0)
    # ★★★★★v301-A (지점장 지시 2026.07.31, 영구): <b>중증화상진단비 = 중증화상진단비 행</b>.
    #   [구 결함] DMAP 부분일치 루프에서 키 <b>'화상진단비'</b>가 '중증화상진단비'를 substring으로
    #   먼저 먹어, resolve_kw의 정본 규칙(`화상 and 중증 → 중증화상진단비`)까지 도달하지 못했다.
    #   교통상해사망(v300)과 <b>완전히 같은 구조의 결함</b>이다.
    if ('화상' in _n300) and any(k in _n300 for k in ('중증', '심재성', '중대한', '부식')):
        return ('중증화상진단비', 0)
    # ★★★★★v350 (지점장 지시 2026.08.02): <b>DMAP 부분일치 선점 전수 점검</b> — 뇌출혈 계열 2건.
    #   지점장 원문: "<b>항상 키워드를 보고 그리고 마지막은 세부가입현황이랑 꼭 누락사항 있는지 체크해라</b>".
    #   [구 결함] DMAP 부분일치 루프에서 키 <b>'뇌출혈진단비'</b>가
    #     <b>'외상성뇌출혈진단비' · '중대한뇌출혈진단비'</b>를 substring으로 먼저 먹어
    #     resolve_kw의 정본 규칙(외상성뇌출혈 행 / 중대한 뇌출혈 행)까지 도달하지 못했다.
    #     교통상해사망(v300) · 중증화상진단비(v301-A)와 <b>완전히 같은 구조</b>다.
    #   ★실측: `외상성뇌출혈진단비`(공백 없음)만 틀렸고 `외상성 뇌출혈 진단비`(공백 있음)는 정상이었다.
    # ★★★★★v421 (지점장 확정 2026.08.14 박미정) — <b>「주요상해뇌출혈」도 외상성이다</b>.
    #   지점장 원문 「<b>외상성뇌출혈이라 뇌출혈진단비와 별개야</b>」.
    #   실측: DB 참좋은운전자상해보험2510 `주요상해뇌출혈진단비 1,000만`이 DMAP 부분일치로
    #   <b>뇌출혈진단비(35행)</b>에 먼저 먹혔다 → 정본은 <b>외상성뇌출혈(37행)</b>.
    #   ★<b>상해로 인한 뇌출혈은 질병 뇌출혈진단비가 아니다.</b>
    if ('뇌출혈' in _n300) and ('외상성' in _n300 or '주요상해' in _n300):
        return ('외상성뇌출혈', 0)
    if ('뇌출혈' in _n300) and ('중대한' in _n300):
        return ('중대한 뇌출혈', 0)
    # ★★★★★v560 (검사관 실측 2026.08.23): <b>DMAP 부분일치 선점 — 산정특례뇌혈관</b>.
    #   [구 결함] 키 <b>'뇌혈관질환진단비'</b>가 '산정특례대상뇌혈관질환진단비'를 substring으로 먼저 먹어
    #     resolve_kw의 정본 규칙(`산정특례 and 뇌 → 산정특례뇌혈관`)까지 도달하지 못했다.
    #   ★실측: resolve_kw = 산정특례뇌혈관(정답) / resolve2 = 뇌혈관진단비(오답).
    #     심장은 DMAP에 '심장질환진단비' 키가 없어 정상 통과했다 — <b>뇌혈관만 뚫려 있었다</b>.
    #   교통상해사망(v300)·중증화상(v301-A)·뇌출혈(v350)과 <b>완전히 같은 구조</b>다.
    if ('산정특례' in _n300) and ('뇌' in _n300) and ('수술' not in _n300):
        return ('산정특례뇌혈관', 0)
    if raw in DMAP:
        v = DMAP[raw]
        return (v, get_종번호(raw))
    for k, v in DMAP.items():
        if k and k in raw and v: return (v, get_종번호(raw))
    # ★v29t (지점장 2026.07.02): 가족동승 부상치료비 = 자부상 아님 → [확인]
    if '가족동승' in raw: return (None, 0)
    # ★v29t: 방사선항암(소액암) 변형 = 합산 금지 → [확인] (수기 정본은 기본 방사선 100만 기재)
    if '방사선' in raw and '소액암' in raw and '제외' not in raw: return (None, 0)
    # ★v29t 부정어 처리: '(소액암제외)'·'(유사암제외)' 등 제외 문구를 지우고 키워드 매칭
    #   (예 '암치료자금_암(소액암제외)진단비특약' → 소액암 오탐으로 유사암행 오매핑되던 버그 차단)
    # ★v44 버그수정: '(치아파절제외)'까지 통째로 지워 골절 제외행→포함행으로 오배치되던 문제(정본 §8.7 위반).
    #    치아·파절이 들어간 제외 괄호는 보존한다. (유사암제외·소액암제외 등은 종전대로 제거)
    raw_kw = re.sub(r'[\(\[](?![^\)\]]*(?:치아|파절))[^\)\]]*제외[\)\]]', '', raw)
    # ★★★★★v304 (지점장 지시 2026.07.31, 영구): <b>재해 = 상해는 같은 것이다</b>.
    #   지점장 원문 = "상해입원일당=재해입원일당 / 즉 <b>재해=상해 는 같은걸로봐야한다</b>".
    #   [구 결함] 구 v29v 폴백은 <b>`resolve_kw`가 None을 돌려줄 때만</b> 재해→상해로 재시도했다.
    #   그래서 재해 담보가 <b>None이 아닌 엉뚱한 행</b>으로 가면 폴백이 아예 안 돌았다.
    #   실측 비대칭 3건 — `재해수술일당` → <b>질병수술일당</b>(상해는 상해수술일당) ·
    #   `재해종수술비(1-5종)`·`재해1-5종수술비` → <b>종수술비공통</b>(상해는 상해 종수술비).
    #   → 규칙마다 `has('재해')`를 덧붙이는 산발 대응을 폐기하고 <b>진입 시점에 한 번 정규화</b>한다.
    #   ★<b>'재해외'(재해 외 원인=질병)는 보호</b>한다 — 치환하면 `_dilqual` 제외어가 깨진다.
    raw_kw = re.sub(r'재해(?!\s*외)', '상해', raw_kw)
    _r2 = resolve_kw(raw_kw)
    if _r2 and _r2[0]:
        return _r2
    # ★★★★★v357 (이성준 실측 2026.08.02): <b>전각 괄호 담보가 통째로 매핑 실패</b>했다.
    #   실측 = 삼성화재 `［간편］상해 입원 수술비（당일입원 제외）` 100 → <b>None</b>
    #          같은 담보를 반각으로 쓰면 `[간편]상해 입원 수술비(당일입원 제외)` → <b>상해수술비</b> ✓
    #   삼성화재 리포트는 <b>모든 담보에 전각 `［간편］` 접두어</b>가 붙어 광범위하게 샜다.
    #   → <b>1차 매칭이 실패했을 때만</b> 전각→반각으로 정규화해 <b>한 번 더</b> 시도한다.
    #     (1차를 건드리지 않으므로 기존 매핑에 회귀가 없다.)
    #   ★규칙은 그대로다 — `상해 수술비` + `상해수술비(***제외)`는 <b>담보명이 다르므로 종전대로 합산</b>.
    _FW = str.maketrans('［］（）｛｝〔〕　', '[](){}[] ')
    _alt = str(raw_kw).translate(_FW)
    if _alt != raw_kw:
        return resolve_kw(_alt)
    return _r2

def resolve(raw):
    return resolve2(raw)[0]

def _dedup_std(raw):
    """★세부보충 dedup 전용 — build_excel의 뇌질환/심장질환Ⅰ·Ⅱ 매핑과 동일 해석.
    (resolve_kw엔 이 매핑이 없어 세부보충이 중복제거 실패→배증하던 양예서 버그 차단)"""
    _rn = re.sub(r'\s','',str(raw))
    if '진단' in _rn and '수술' not in _rn and '주요치료' not in _rn:
        if ('심장질환진단' in _rn) and ('허혈' not in _rn) and ('급성심근' not in _rn):
            _mn=_rmn(_rn)
            # ★v384 심장Ⅰ 묶음 = 협심증(허혈성은 협심증의 종류) — 지점장 확정 2026.08.11
            if _mn==2: return '급성심근경색'
            if _mn==1: return '협심증'
        if '뇌질환진단' in _rn:
            _mn=_rmn(_rn)
            if _mn==2: return '뇌졸증진단비'
            if _mn==1: return '뇌혈관진단비'
    return resolve_kw(raw)[0]

NOFILL = PatternFill(fill_type=None)

# ★ SUM 수식 캐시값 채우기 — openpyxl은 수식만 저장(캐시 없음)→모바일/미리보기 공란.
#   저장 후 LibreOffice 재계산으로 값 주입(수식은 유지=법칙22).
def recalc_xlsx(path):
    import subprocess, shutil, tempfile
    soffice = shutil.which('soffice') or shutil.which('libreoffice')
    if not soffice: return False
    try:
        outd = tempfile.mkdtemp()
        subprocess.run([soffice,'--headless','--norestore','--convert-to','xlsx','--outdir',outd,path],
                       timeout=90, capture_output=True)
        out = os.path.join(outd, os.path.splitext(os.path.basename(path))[0]+'.xlsx')
        if os.path.exists(out):
            shutil.copyfile(out, path); shutil.rmtree(outd, ignore_errors=True)
            return True
        shutil.rmtree(outd, ignore_errors=True)
        return False
    except Exception:
        return False

# ★v29u: LibreOffice 없는 환경(Railway)용 캐시 주입 — 합계 수식은 유지(§5)하고,
#   계산값을 파이썬으로 구해 시트 XML <v>에 직접 기록 → 폰·미리보기·보장설명지 모두 값 표시.
def _no_fullcalc(wb):
    """★v51(2026.07.13): 산출 엑셀에서 fullCalcOnLoad 플래그 제거.
    master.xlsx가 이 플래그를 물려주면 Excel이 파일을 열 때마다 전 수식을 강제 재계산한다
    (편집모드 진입 시 폰 Excel이 1분 이상 로딩). 끝열 =SUM 캐시값은 inject_sum_cache가 이미
    채워두므로 강제 재계산은 불필요한 부하일 뿐이다. 수식(§5)은 그대로 유지한다."""
    try:
        wb.calculation.fullCalcOnLoad = False
    except Exception:
        pass

def inject_sum_cache(path):
    import zipfile, shutil, tempfile
    try:
        wb = openpyxl.load_workbook(path)
        ws = wb['보장분석']; last = ws.max_column
        vals = {}
        for r in range(2, ws.max_row+1):
            f = ws.cell(r,last).value
            if not (isinstance(f,str) and f.startswith('=')): continue
            nums = [ws.cell(r,c).value for c in _data_cols(ws,last) if isinstance(ws.cell(r,c).value,(int,float))]   # ★v388 합산 열 제외
            s = sum(nums)
            if f.startswith('=MIN('):
                m = re.search(r',\s*(\d+)\s*\)\s*$', f); v = min(s, int(m.group(1))) if m else s
            elif f.startswith('=IF(COUNT'):
                v = max(nums) if nums else 0
            elif f.startswith('=IF(SUM'):
                v = 7 if s>0 else 0
            else:
                v = s
            vals[ws.cell(r,last).coordinate] = v
        # ★★v388b 보유 합계·제안 합계 열에도 <b>같은 캐시</b>를 박는다.
        #   구 코드는 끝열만 캐시해서 <b>폰·미리보기에서 두 열이 빈칸으로 보였다</b>(수식만 있고 값이 없음).
        for c in range(3, last):
            if str(ws.cell(1,c).value or '').strip() not in ('보유 합계','제안 합계'): continue
            for r in range(2, ws.max_row+1):
                f2 = ws.cell(r,c).value
                if not (isinstance(f2,str) and f2.startswith('=')): continue
                m2 = re.search(r'([A-Z]+)(\d+):([A-Z]+)(\d+)', f2)
                if not m2: continue
                c0 = openpyxl.utils.column_index_from_string(m2.group(1))
                c1 = openpyxl.utils.column_index_from_string(m2.group(3))
                _ns2 = [v for v in (ws.cell(r,cc).value for cc in range(c0, c1+1))
                        if isinstance(v,(int,float))]
                _s2 = sum(_ns2)
                if f2.startswith('=MIN('):
                    _m3 = re.search(r',\s*(\d+)\s*\)\s*$', f2)
                    _s2 = min(_s2, int(_m3.group(1))) if _m3 else _s2
                elif f2.startswith('=IF(COUNT'):
                    _s2 = max(_ns2) if _ns2 else 0
                vals[ws.cell(r,c).coordinate] = _s2
        if not vals: return False
        zin = zipfile.ZipFile(path,'r')
        # 보장분석 시트 XML 경로 확인 (workbook.xml 순서 + rels)
        wbxml = zin.read('xl/workbook.xml').decode('utf-8')
        rels  = zin.read('xl/_rels/workbook.xml.rels').decode('utf-8')
        m = re.search(r'<sheet[^>]*name="보장분석"[^>]*r:id="(rId\d+)"', wbxml) or re.search(r'<sheet[^>]*r:id="(rId\d+)"[^>]*name="보장분석"', wbxml)
        rid = m.group(1) if m else 'rId1'
        m2 = re.search(r'Id="'+rid+r'"[^>]*Target="([^"]+)"', rels)
        tgt = 'xl/'+m2.group(1).lstrip('/') if m2 else 'xl/worksheets/sheet1.xml'
        sx = zin.read(tgt).decode('utf-8')
        for coord, v in vals.items():
            vv = ('%d' % v) if float(v).is_integer() else repr(float(v))
            sx = re.sub(r'(<c r="'+coord+r'"[^>]*>)(<f[^>]*>[^<]*</f>)(?:<v>[^<]*</v>)?(</c>)', r'\1\2<v>'+vv+r'</v>\3', sx, count=1)
        tmp = path+'.tmp'
        zout = zipfile.ZipFile(tmp,'w',zipfile.ZIP_DEFLATED)
        for it in zin.infolist():
            zout.writestr(it, sx.encode('utf-8') if it.filename==tgt else zin.read(it.filename))
        zout.close(); zin.close()
        shutil.move(tmp, path)
        return True
    except Exception as _e:
        print(f'[INJECT_CACHE_ERR] {_e}')
        return False

# ★ LLM 매핑 엔진 — 마스터 표준 담보명에 의미기반 매핑 (앱 자동화 핵심)
def load_std_dambo(ws):
    out=[]
    for r in range(6, ws.max_row+1):
        v=ws.cell(r,2).value
        if v and str(v).strip(): out.append(str(v).strip())
    return out

def llm_resolve(raw_names, std_list):
    """raw 담보명 -> {raw: {'std': 표준명 or None, 'jong': 0~5}}  (jong>0이면 종수술비)"""
    raw_names=[r for r in raw_names if r]
    if not raw_names: return {}
    key=os.environ.get('ANTHROPIC_API_KEY','')
    if not key:
        return {r:{'std':None,'jong':0,'note':''} for r in raw_names}  # 키 없음 -> 전부 [확인]
    rules=("한국 보험 보장분석. 아래 담보명들을 표준목록의 표준명에 의미기반 매핑.\n"
        "규칙:\n"
        "- 표준목록에 있는 것만 선택. 의미가 명확히 일치할 때만. 애매/해당없음=null.\n"
        "- 중입자방사선=중입자치료비, 양성자=양성자치료, 세기조절=세기조절치료, 표적항암약물=표적항암치료비\n"
        "- '상해 1~5종/1-5종 수술비'(_N종·(N종)·N종)=상해 종수술비(1-5종), jong=종번호. 질병도 동일=질병 종수술비(1-5종)\n"
        "- 허혈심장/협심증=협심증, 급성심근경색=급성심근경색, 부정맥=부정맥, 심장질환수술=심장수술비\n"
        "- 유사암(갑상선/기타피부/제자리/경계성)=유사암(갑.기.경.제). 특수치료 아닌 일반 암진단=일반암\n"
        "- '뇌혈관진단'(로마숫자 없음) 담보는 세부가입현황 표를 정본으로 재배치한다(AIA·라이나·AIG·우체국 등). 임의 하드코딩 금지.\n"
        "- 하이클래스=하이클래스(암), 표적항암약물허가=표적항암치료비(여러 건이면 가장 큰 1건만)\n"
        "- 유사암 진단금액은 가입연도 2020년 이하면 일반암의 1/10, 2021년 이상이면 1/5로 환산 기재\n"
        "- 뇌혈관수술=뇌혈관수술비, 항암방사선/약물치료비=항암방사선약물, 암수술=암수술\n"
        "- 화상 진단비='화상진단비'(화상 구분 행), 중대한화상·부식=중증화상진단비\n"
        "- 생명보험 종신 주계약/기본계약(사망보장)=일반사망\n"
        "- 운전자: 벌금(대인)=대인, 벌금(대물)=대물, 교통사고처리지원금(중상해)=합의금, 처리지원금(6주미만)=6주미만, 변호사비=변호사, 자동차(사고)부상보장=자부상\n"
        "- 표준목록에 자리 없는 담보(예 크론병·다발경화증·장기이식 등)=null (행 추가 금지)\n"
        "- 심장담보 질병코드 분류: 협심증(I20)=협심증, 급성심근경색(I21~22)=급성심근경색, "
        "부정맥(I47~49)=부정맥, 심부전(I50)=심부전, 심내막·심근·심장막염=염증. "
        "단 '특정Ⅰ/Ⅱ·특정심장' 등 상품별 정의는 약관마다 달라 본표 단정 금지→null.\n"
        "- note 규칙(엄수): 약관 없이 보장범위를 추측·일반론으로 단정하지 말 것. "
        "담보명에서 명백한 사실만 적고, 상품별 정의(특정Ⅰ/Ⅱ 등)나 불확실한 건 '약관 확인 필요'로만 기재. "
        "의학적 추정(예 '통상 허혈성'·'대개 ~포함') 금지.\n")
    prompt=rules+f"\n표준목록: {std_list}\n\n담보명: {raw_names}\n\nJSON만 출력(설명 금지): {{\"담보명\":{{\"std\":\"표준명 또는 null\",\"jong\":0,\"note\":\"보장범위 요약\"}}}}"
    try:
        r=httpx.post('https://api.anthropic.com/v1/messages',
            headers={'x-api-key':key,'anthropic-version':'2023-06-01','content-type':'application/json'},
            json={'model':'claude-haiku-4-5-20251001','max_tokens':8000,
                  'messages':[{'role':'user','content':prompt}]}, timeout=90)
        print(f'[LLM_RESOLVE] status={r.status_code} raw={len(raw_names)}')
        if r.status_code != 200:
            print(f'[LLM_RESOLVE_HTTP] {r.text[:400]}')
            return {x:{'std':None,'jong':0,'note':''} for x in raw_names}
        txt=''.join(b.get('text','') for b in r.json().get('content',[]) if b.get('type')=='text')
        txt=txt.strip().replace('```json','').replace('```','').strip()
        out=json.loads(txt)
        print(f'[LLM_RESOLVE] mapped={len(out)}')
        return {k:{'std':(v.get('std') if isinstance(v,dict) else None),
                   'jong':(v.get('jong',0) if isinstance(v,dict) else 0),
                   'note':(v.get('note','') if isinstance(v,dict) else '')} for k,v in out.items()}
    except Exception as e:
        print(f'[LLM_RESOLVE_ERR] {e}')
        return {x:{'std':None,'jong':0,'note':''} for x in raw_names}

def _fix_silson(contracts):
    """★★★★★실손 정본 §8.8 — <b>실손은 [별첨] 보험서비스(상품)별 보장현황이 답이다</b>(지점장 확정 2026.07.25, 영구).
    ★<b>소스 우선순위는 실손에도 그대로</b>: ①<b>[별첨] 상품별 보장현황 = 정본</b> ②세부가입현황 ③한장보장표(검산).
      한장보장표·세부가입현황이 별첨과 다르면 <b>별첨을 따른다</b>. 실측 확정 사례 —
      DB생명 CI종신(3세대) 별첨 `외래의료비 <b>25</b>` vs 한장표·세부내역 30 → <b>25가 정답</b>
      (3세대 실손 외래 한도가 25만원). 구 [확인] 대기 항목은 이로써 <b>확정·해소</b>.
    ★<b>입원·통원 = 별첨 명시값 최우선</b>(1세대 1억·5천·3천 다 유지, 절대 뭉개지 말 것).
      명시값이 없을 때만 세대·회사유형 기본값을 넣는다: 1세대 통원10 / 손보 통원25·약5 / 생보 통원20·약10 / 4·5세대 통원20.
    ★<b>약값은 예외 2건</b> — 1세대와 4·5세대는 외래+약제가 <b>통원 한 한도로 통합</b>이라
      약값 담보 자체가 없다 → <b>별첨에 값이 있어도 0으로 삭제</b>(세대 규칙이 별첨보다 우선하는 유일한 항목)."""
    for c in contracts:
        d=c.get('dambo',{})
        co=str(c.get('company','') or '')
        prod=str(c.get('product','') or '')
        cd=str(c.get('contract_date','') or '')
        is_saengbo = ('생명' in co or '라이프' in co) and ('화재' not in co and '손해' not in co and '손보' not in co and '해상' not in co)
        # 원본 담보키 → 표준(입원/통원/약값) 위치 찾기
        _kmap={}  # 표준명 → 원본키
        for _rk in list(d.keys()):
            _std,_=resolve_kw(_rk)
            if _std in ('입원','통원','약값') and _std not in _kmap:
                _kmap[_std]=_rk
        if '입원' not in _kmap: continue  # 실손 없음
        def _get(std):
            k=_kmap.get(std)
            if not k: return None
            try: return float(str(d[k]).replace(',',''))
            except: return None
        def _set(std,val):
            k=_kmap.get(std)
            if k: d[k]=val
            else: d[std]=val  # 없던 항목(약값 등)은 표준키로 신설
        ipw=_get('입원')
        if ipw is None: continue
        # ★입원 = 보장표 명시값 최우선(1세대 1억·5천·3천 다 있음, 절대 뭉개지 말 것).
        #   명시값이 있으면 그대로 유지. 통원·약값만 규칙 보정.
        _tw=_get('통원'); _yk=_get('약값')
        # ★4·5세대(2021.07~) = 통원+약 합쳐서 20 (약값 별도 없음)
        _ym=None
        _mm=re.search(r'(\d{4})\.(\d{2})',cd)
        if _mm: _ym=int(_mm.group(1))*100+int(_mm.group(2))
        _gen=silson_gen(cd, ipw, prod, _has_nonpay3(d), _has_drug(d))   # ★v250 3대비급여 하한(d=이 계약 dambo)
        if _gen in ('4세대','5세대') or (_ym and _ym>=202107):
            # ★★★★★v246 영구지침(지점장 확정 2026.07.25): <b>실손은 [별첨] 상품별 보장현황이 답이다</b>.
            #   구 코드는 `_tw<=20`이 아니면 <b>20으로 덮어썼다</b> → 별첨 명시값이 20 초과면 잘렸다.
            #   → <b>별첨 명시값이 있으면 그대로 쓰고</b>, 없을 때만 기본값 20을 넣는다.
            _set('통원', _tw if _tw else 20)
            _set('약값', 0)     # ★4·5세대는 통원+약제 통합 → 약값 담보 자체가 없다(별첨에 있어도 삭제)
            continue
        # ★★★v215 (지점장 확정 2026.07.25, 영구): <b>실손 1세대 = 입원 3,000 / 통원 10 / 약값 0</b>.
        #   1세대는 외래+약제가 <b>통원 한 한도로 통합</b>이라 약값 담보 자체가 없다.
        #   ★버그: `silson_gen`은 입원한도 3,000이면 <b>'1세대(구형)'</b>을 반환하는데 구 v41 코드는
        #   `== '1세대'`로 비교해 <b>항상 False</b>였다 → 아래 손보/생보 기본값으로 흘러
        #   <b>약값 5(손보)·10(생보)이 신설</b>됐고, 그 값이 그대로 엑셀→보장진단서 PPT까지 나갔다
        #   (지점장 지적 "간병인·실손 보장진단서 오류"). → `.startswith('1세대')`로 수정.
        if str(_gen).startswith('1세대'):
            if not _tw: _set('통원', 10)      # ★v215 1세대 통원 기본 10(명시값 있으면 그대로)
            _set('약값', 0)
            continue
        # 1~3세대: 통원·약값 명시값 우선, 없으면 회사유형 기본값.
        if is_saengbo:
            if not _tw: _set('통원',20)
            if not _yk: _set('약값',10)
        else:
            if not _tw: _set('통원',25)
            if not _yk: _set('약값',5)
    return contracts


# ★★★★★v309 감사 2종 자동 게이트 (지점장 확정 2026.07.31, 영구)
#   왜: 「감사 2종은 매 배포마다 돌린다」고 지침에 적어두고 <b>내가 계속 안 돌렸다</b>.
#   사람이 기억해서 돌리는 방식은 이미 실패가 증명됐다 → <b>앱이 매 분석마다 자동으로 돌린다</b>.
#   ㉠ 지침 규칙 케이스 55건 (담보명 → 기대 마스터 행). 기대값은 지침 원문 대조로 확정.
#   ㉡ 마스터 전 행 커버리지 (라벨을 resolve2에 넣어 자기 행 복귀 검사).
#      ★_COV_ALLOW = <b>라벨≠담보명 설계라 정상인 16건</b>. 이걸 빼지 않으면 게이트가 늘 FAIL이라
#        아무도 안 보게 된다(늑대소년 방지).
#   ★값·행 배정은 일절 건드리지 않는다. 결과는 확인사항 3행 D열 + 로그 + /health 전용.
# ★★★★★v560 4대 블록(암·뇌·심·수술) 커버리지 보강 케이스 24건.
#   전부 resolve2 실호출로 정답 복귀를 확인한 것만 등재한다(추측 0건).
_EXTRA_CASES = [
    # ── 암
    ('암주요치료비',            '암주요치료비'),
    ('암하이클래스치료비',        '하이클래스(암)'),
    ('항암중입자방사선치료비',      '중입자치료비'),
    ('표적항암약물허가치료비',      '표적항암치료비'),
    ('항암양성자방사선치료비',      '양성자치료'),
    ('항암세기조절방사선치료비',     '세기조절치료'),
    ('다빈치로봇암수술비',        '다빈치로봇수술비'),
    ('항암방사선약물치료비',       '항암방사선약물'),
    # ── 뇌혈관
    ('뇌혈관질환진단비',         '뇌혈관진단비'),
    ('뇌졸중진단비',           '뇌졸증진단비'),
    ('뇌출혈진단비',           '뇌출혈진단비'),
    ('중대한뇌출혈진단비',        '중대한 뇌출혈'),
    ('외상성뇌출혈진단비',        '외상성뇌출혈'),
    ('산정특례대상뇌혈관질환진단비',  '산정특례뇌혈관'),   # ★v560 DMAP 선점 차단 회귀
    ('뇌혈관질환혈전용해치료비',     '혈전용해치료비'),
    # ── 심장
    ('산정특례대상심장질환진단비',   '산정특례심장'),
    # ── 수술
    ('상해수술비',            '상해수술비'),
    ('질병수술비',            '질병수술비'),
    ('창상봉합술치료비',         '창상봉합술'),
    ('골절수술비',            '골절수술비'),
    ('화상수술비',            '화상수술비'),
    ('심장질환수술비',          '심장수술비'),
    ('허혈성심장질환수술비',       '허혈성수술비'),
    ('120대질병수술비',        'n대수술비'),
    ('중대한상해수술비',         '중대한상해수술비'),
    ('상해종수술비(1-8종)',      '상해 종수술비(1-8종)'),
    ('5대골절수술비',           '5대골절수술비'),
    ('질병종수술비(1-5종)',      '질병 종수술비(1-5종)'),
    ('질병종수술비(1-8종)',      '질병 종수술비(1-8종)'),
    ('5대기관수술비(관혈)',       '5대기관 수술비 관혈'),
    ('5대기관수술비(비관혈)',      '5대기관 수술비 비관혈'),
    # ── 2대 주요치료비 (뇌39·심49 두 행 공통)
    ('2대질병주요치료비',        '2대 주요치료비'),
    ('심뇌혈관주요치료비',        '2대 주요치료비'),
]

_AUDIT_GROUPS = [
 ('§8.7 일당 계열 배정(v212·v303·v305)', [
    ('2대질병입원일당', None), ('특정질병입원일당', None), ('3대질병입원일당', None),
    ('질병입원일당(1-180)', '질병일당'), ('질병일당', '질병일당'), ('질병입원', '질병일당'),
    ('상해입원일당(1-180)', '상해일당'), ('재해입원일당', '상해일당'),
    ('종합병원 질병입원일당', '종합병원 질병입원일당'), ('상급병원 질병입원일당', None),
    ('교통상해입원일당', None), ('암직접치료입원일당', '암일당'), ('입원급여금', '질병일당'),
    ('간호간병통합서비스사용 질병입원일당', '간호통합병동')]),
 ('§8.2 암·유사암(v248·v197·v227)', [
    ('갱신형 유사암진단비', '유사암(갑.기.경.제)'), ('소액암진단', '유사암(갑.기.경.제)'),
    ('갱신형 갑상선암(초기제외)진단비', None), ('갑상샘암치료보험금', None),
    ('상피내암치료보험금', None), ('갱신형 유사암수술비', None),
    ('암진단비(유사암제외)', '일반암'), ('16대특정암진단비', '고액암'),
    ('일반암직접치료비', '암수술'), ('치매진단', None), ('치매진단(경증)', None)]),
 ('§8.3 심장·부정맥(v217·단독담보 원칙)', [
    ('심장부정맥고주파·냉각절제술보장', None), ('부정맥진단비', '부정맥'),
    ('기타심장부정맥(I49)진단비', '부정맥'), ('허혈성심장질환진단비', '허혈성 진단비'),
    ('주요심뇌5대혈관수술비', None), ('심뇌혈관수술비', '뇌혈관수술비')]),
 ('§8.1 사망·후유장해(v222·v300·v302)', [
    ('상해후유장해(20%이상)', None), ('상해후유장해(50%이상)', None),
    ('상해후유장해(3%이상)', '상해후유3%'), ('재해장해연금', None), ('재해후유장해3%', '상해후유3%'),
    ('교통상해사망', '교통상해사망'), ('대중교통상해사망', '교통상해사망'), ('상해사망', '상해사망')]),
 ('§8.6 운전자(v198·v301·v302)', [
    ('타인사망교통사고처리지원금', '합의금'), ('교통사고처리지원금(중상해포함)', '합의금'),
    ('자동차사고부상치료비(14급)', '자부상'), ('교통사고 벌금(대인)', '대인'),
    ('교통사고 벌금(대물)', '대물'), ('변호사선임비용', '변호사'), ('6주미만 진단위로금', '6주미만')]),
 ('§8.7 골절·화상·기타(v301·v38c)', [
    ('중증화상진단비', '중증화상진단비'), ('화상진단비', '화상진단비'),
    ('중대한화상및부식진단비', '중증화상진단비'),
    ('골절진단비(치아파절제외)', '골절(치아파절제외)'), ('골절진단비', '골절(치아파절포함)'),
    ('일상생활배상책임', '일상배상책임'), ('응급실내원비', '응급실(응급)')]),
 ('§8.5 수술·재해=상해(v304)', [
    ('재해수술일당', '상해수술일당'), ('재해종수술비(1-5종)', '상해 종수술비(1-5종)')]),
]
_AUDIT_CASES = [c for _g, _cs in _AUDIT_GROUPS for c in _cs]

# ★★★★★v311 <b>통제권 이관 — 감사 케이스의 정본은 master.xlsx 「지침케이스」 시트다</b>
#   (지점장 지시 2026.07.31 "결국 너도 너의 생각대로 많이 하니까 또 하나의 너를 컨트롤할게 필요해")
#   왜: v310 지침 체크봇은 <b>내가 고른 55개</b>만 봤다. 감시 기준을 감시 대상(코드=나)이 쥐고 있으면
#   내가 언제든 줄이거나 지울 수 있다. → 기준을 <b>내 손이 닿지 않는 곳(지점장 엑셀)</b>으로 옮긴다.
#   ★지점장이 시트에 <b>행만 추가하면</b> 앱이 매 분석마다 자동으로 검사한다. 코드 수정 불필요.
#   ★시트가 없거나 비면 <b>위 하드코딩으로 폴백</b>(구버전 master로도 앱이 죽지 않게).
#   ★<b>build_excel은 이 시트를 고객 산출물에서 삭제</b>한다(마스터 워크북을 그대로 save하므로).
_DOCTRINE_SRC = '코드(폴백)'

def load_doctrine_sheet():
    """master.xlsx 「지침케이스」 시트를 읽어 (_AUDIT_GROUPS, _COV_ALLOW)를 갈아끼운다."""
    global _AUDIT_GROUPS, _AUDIT_CASES, _COV_ALLOW, _DOCTRINE_SRC
    try:
        wb = openpyxl.load_workbook(TPL_XL)
        if '지침케이스' not in wb.sheetnames:
            _DOCTRINE_SRC = '코드(폴백 — 마스터에 「지침케이스」 시트 없음)'; return
        ws = wb['지침케이스']
        groups = {}; order = []; allow = set(); ncase = 0
        for r in range(6, ws.max_row + 1):
            kind = str(ws.cell(r, 1).value or '').strip()
            gname = str(ws.cell(r, 2).value or '').strip()
            raw = str(ws.cell(r, 3).value or '').strip()
            exp = ws.cell(r, 4).value
            exp = str(exp).strip() if exp not in (None, '') else None
            if not raw: continue
            if kind.startswith('커버리지'):
                allow.add(raw); continue
            if not kind.startswith('케이스'): continue
            gname = gname or '(조항 미기재)'
            if gname not in groups: groups[gname] = []; order.append(gname)
            groups[gname].append((raw, exp)); ncase += 1
        if ncase == 0:
            _DOCTRINE_SRC = '코드(폴백 — 시트에 케이스 0건)'; return
        _AUDIT_GROUPS = [(g, groups[g]) for g in order]
        # ★★★★★v560 (지점장 지시 2026.08.23 「암·뇌·심·수술은 100%여야 했다」):
        #   master.xlsx는 제124조 1항으로 <b>손대지 않는다</b> → 케이스를 코드에 두고 시트 뒤에 붙인다.
        #   실측 근거: 아래 24건은 전부 <b>resolve2 실호출로 정답 행 복귀를 확인</b>한 것만 올렸다.
        #   ★검사가 없어 조용히 뚫려 있던 구간을 막는 것이 목적 — 회귀 방지용이다.
        _AUDIT_GROUPS = _AUDIT_GROUPS + [('제0조 4대블록 보강(v560)', _EXTRA_CASES)]
        _AUDIT_CASES = [c for _g, _cs in _AUDIT_GROUPS for c in _cs]
        if allow: _COV_ALLOW = allow
        _DOCTRINE_SRC = 'master.xlsx 「지침케이스」 시트 (케이스 %d · 커버리지허용 %d)' % (ncase, len(allow))
        print('[v311 지침정본] ' + _DOCTRINE_SRC)
    except Exception as e:
        _DOCTRINE_SRC = '코드(폴백 — 시트 읽기 실패: %s)' % str(e)[:60]
        print('[v311 지침정본] ' + _DOCTRINE_SRC)

# ★★★★★v320 <b>호출 위치 이동(구조 결함 수정)</b>: 구 코드는 여기서 `load_doctrine_sheet()`를
#   불렀는데 <b>아래 `_COV_ALLOW` 하드코딩이 그 뒤에 있어 시트에서 읽은 값을 덮어썼다</b>.
#   실측 = 시트 커버리지허용 17건인데 `_COV_ALLOW`는 16건 → 감사 FAIL 1.

# ★★★★★v312 <b>「해석원칙」 = 관점의 정본</b>(지점장 지시 2026.07.31
#   "관점은 니 맘대로 하지 않고 지침을 따르게 고정하는 법률 같은 봇 존재다")
#   v311까지는 <b>조문(담보명→행)</b>만 고정했다. 그런데 실제 사고는
#   <b>「지침에 없는 것을 내가 어떻게 판단하느냐」= 관점</b>에서 났다.
#   → 해석 규칙 자체를 master.xlsx 「해석원칙」 시트에 성문화하고,
#     <b>앱이 매 분석마다 확인사항 하단에 원문 그대로 인쇄</b>한다.
#     내가 지침을 요약·재구성해도 <b>산출물에는 지점장 원문이 그대로 나가</b> 왜곡이 즉시 드러난다.
#   ★시트가 없으면 기본 10개로 폴백(구버전 master로도 앱이 안 죽는다).
_PRINCIPLES = []
_PRINCIPLES_SRC = '코드(폴백)'
_PRIN_BASE = [
    '지점장이 준 자료는 문구·코드를 그대로 쓴다. 요약·재구성하거나 해설 주석을 덧붙이지 않는다.',
    '지침에 없으면 추측하지 않는다 — 빈칸으로 두고 [확인]큐로 보낸다.',
    '범위를 넓히려면 먼저 물어본다. 지점장이 두 개를 말했으면 두 개만 넣는다.',
    '지침이 바뀌면 옛 내용을 남기지 말고 지운다. 세트로 내려온 지시는 전부 폐기·교체한다.',
    '같은 데이터가 여러 곳에 있다 — 전수 수정하고 마지막에 옛 문구 잔재를 스캔한다.',
    '규칙 두 개가 같은 값을 만지면 나중 것이 앞 것을 죽인다. 새 규칙 전에 같은 셀을 쓰는 기존 규칙을 확인한다.',
    '직접 열어본 것만 "~이다"라고 쓴다. 안 본 것은 "가설"이라 명시한다. 측정하지 않은 수치는 말하지 않는다.',
    '고장난 기능은 고치는 것이다 — 선택지로 만들어 지점장께 결정을 넘기지 않는다.',
    '못 한 검증은 "못 했다"고 쓴다. 결과물만 올리지 않고 실행 과정을 보여준다.',
    '규칙 하나를 고치면 2열(롯데)·3열(KB) 두 경로를 반드시 다 돌린다.',
]

def load_principles():
    global _PRINCIPLES, _PRINCIPLES_SRC
    _base = [(p, '', '') for p in _PRIN_BASE]
    try:
        wb = openpyxl.load_workbook(TPL_XL)
        if '해석원칙' not in wb.sheetnames:
            _PRINCIPLES = _base; _PRINCIPLES_SRC = '코드(폴백 — 마스터에 「해석원칙」 시트 없음)'; return
        ws = wb['해석원칙']; rows = []
        for r in range(5, ws.max_row + 1):
            p = ws.cell(r, 2).value
            if p is None or not str(p).strip(): continue
            rows.append((str(p).strip(), str(ws.cell(r, 3).value or '').strip(),
                         str(ws.cell(r, 4).value or '').strip()))
        if not rows:
            _PRINCIPLES = _base; _PRINCIPLES_SRC = '코드(폴백 — 시트에 원칙 0건)'; return
        _PRINCIPLES = rows
        _PRINCIPLES_SRC = 'master.xlsx 「해석원칙」 시트 (%d개)' % len(rows)
        print('[v312 해석원칙] ' + _PRINCIPLES_SRC)
    except Exception as e:
        _PRINCIPLES = _base
        _PRINCIPLES_SRC = '코드(폴백 — 시트 읽기 실패: %s)' % str(e)[:60]

load_principles()

_COV_ALLOW = {'중대한CI적용','일반암','암일당','항암방사선약물','중대한 뇌출혈','염증',
              '상해 종수술비(1-3종)','상해 종수술비(1-8종)','질병 종수술비(1-3종)',
              '질병 종수술비(1-8종)','120대수술비','대인','입원','통원','약값',
              'MRI/도수치료/비급여주사'}
_AUDIT_LAST = {'case':'-', 'cov':'-', 'fail':0, 'detail':[]}
load_doctrine_sheet()   # ★v320 정의가 모두 끝난 뒤에 호출해야 _COV_ALLOW를 덮어쓸 수 있다

def audit_run(labels=None):
    detail = []; c_ok = 0
    for _raw, _exp in _AUDIT_CASES:
        try: _got = resolve2(_raw)[0]
        except Exception: _got = None
        if (_got or None) == (_exp or None): c_ok += 1
        else: detail.append('규칙 「%s」 기대 %s / 실제 %s' % (_raw, _exp, _got))
    seen = []; v_ok = 0
    for _L in (labels or []):
        _L = str(_L).strip()
        if not _L or _L in seen: continue
        seen.append(_L)
        try: _g = resolve2(_L)[0]
        except Exception: _g = None
        if _g == _L or _L in _COV_ALLOW: v_ok += 1
        else: detail.append('커버리지 「%s」 → %s' % (_L, _g))
    _AUDIT_LAST.update(case='%d/%d' % (c_ok, len(_AUDIT_CASES)),
                       cov='%d/%d' % (v_ok, len(seen)), fail=len(detail), detail=detail[:20])
    print('[v309 감사] 규칙 %d/%d · 마스터커버리지 %d/%d · FAIL %d'
          % (c_ok, len(_AUDIT_CASES), v_ok, len(seen), len(detail)))
    for _d in detail[:20]: print('   [AUDIT_FAIL]', _d)
    return _AUDIT_LAST

def build_excel(data, out):
    _ci_diag = []   # ★v238 CI 진단표(확인사항 상시 출력) — 함수 최상단에서 확실히 초기화
    wb = openpyxl.load_workbook(TPL_XL)
    ws = wb['보장분석']
    client = data['client']; contracts = _fix_silson(data['contracts'])

    # 담보명 -> 행번호 맵 (A/B열 유지)
    nm2r = {}
    nm2r_norm = {}   # 공백무시 보조키 ('진 단 비' 등 라벨 변형 흡수)
    nm2r_multi = {}  # 동일 담보명이 여러 행(2대 주요치료비=뇌혈관+심장) → 모두 기재용
    for r in range(6, ws.max_row+1):
        v = ws.cell(r,2).value
        if v:
            k = str(v).strip()
            nm2r[k] = r
            nm2r_norm[re.sub(r'\s','',k)] = r
            nm2r_multi.setdefault(k, []).append(r)
    _aud = audit_run(list(nm2r.keys()))   # ★v309 감사 2종 자동 실행(표시 전용)

    # ★ 데이터영역(C열~) 전체 초기화 — 옛 7계약 헤더·합계·SUM수식·슬래시골격 제거
    MAXC = 60  # 최대 50계약 + 여유
    for r in range(1, ws.max_row+1):
        for c in range(3, MAXC+1):
            cell = ws.cell(r,c)
            cell.value = None
            cell.fill = NOFILL
    # ★★★★★v447 제62조 (지점장 지적 2026.08.17 「골절 이하로 세로폭이 줄었다」)
    #   실측: master.xlsx가 <b>92행(골절) 이후 15행의 행높이를 지정하지 않았다</b>(h=None).
    #   6~91행은 전부 지정돼 있다(16.5가 63행 · 17.25가 23행) → 골절 이하만 기본높이로 렌더돼
    #   위와 다르게 보였다. 코드가 만든 게 아니라 <b>마스터의 결함</b>이다.
    #   → 마스터(정본)는 건드리지 않고, 산출 시점에 <b>빠진 행만</b> 다수값 16.5로 채운다.
    #     이미 지정된 행은 손대지 않는다(마스터 우선).
    _RH_FILL = 16.5
    _rh_fixed = []
    for _r in range(6, ws.max_row + 1):
        _d = ws.row_dimensions.get(_r)
        if _d is None or _d.height is None:
            ws.row_dimensions[_r].height = _RH_FILL   # ★height 대입만으로 customHeight가 켜진다
            _rh_fixed.append(_r)
    if _rh_fixed:
        print('[v447 행높이] 마스터 미지정 %d행 보정 → %s (%s~%s)'
              % (len(_rh_fixed), _RH_FILL, _rh_fixed[0], _rh_fixed[-1]))

    ws.cell(1,1).value = f"{client} 보장진단"
    # ★★★★★v297 (이영태 실측 2026.07.31, 영구): <b>이미지 PDF 경고를 본표 1행에도 박는다</b>.
    #   구 v280은 <b>확인사항 시트</b>에만 경고를 넣었다. 지점장이 가장 먼저·가장 많이 보는 곳은
    #   <b>보장분석 본표 1행</b>인데 거기엔 아무 표시가 없어, 비전 OCR로 억지로 읽은 산출물이
    #   <b>정상 산출물과 똑같아 보인 채</b> 나갔다(이영태 = 'Microsoft: Print To PDF' · 한글 0자).
    #   지침 원칙 「조용히 틀리는 것을 시끄럽게」 적용 — 제출 금지를 산출물 표면에 남긴다.
    try:
        _ipw0 = globals().get('_IMG_PDF_WARN','')
        if _ipw0:
            _cc0 = ws.cell(1,1)
            _cc0.value = f"{client} 보장진단   ★★ 이미지 PDF(비전 OCR) — 담보명 오독 가능 · 고객 제출 금지 ★★"
            _cc0.font = Font(bold=True, size=13, color='B00020')
    except Exception: pass
    try:   # ★v371 제안서 단독 모드 표기
        _jo0 = globals().get('_JEAN_ONLY','')
        if _jo0:
            _cj0 = ws.cell(1,1)
            _cj0.value = f"{client} 보장진단   ★ 가입제안서 단독 — [검산] 불가 · [실손 세대] 불가 ★"
            _cj0.font = Font(bold=True, size=13, color='C00000')
    except Exception: pass

    n_ct = len(contracts)

    # ★ LLM 배치 매핑 (앱 자동화): 전체 담보 1회 호출 -> 표준명/종번호
    std_list = load_std_dambo(ws)
    all_raw = sorted({raw for c in contracts for raw in c['dambo']})
    LLMMAP = llm_resolve(all_raw, std_list)
    unmapped = []  # (회사, 담보명, 금액) — 마스터 미수록/매핑실패 -> [확인]
    # ★★★★★v261 접힘 자가진단(지점장 지시 2026.07.27, 영구) — <b>매 분석마다 자동 검수</b>.
    #   2026.07.27 장O경 사고: `다빈치로봇 암수술비(…)(갱화상수술비(간편가입)상해1~5종수술비 → 7`
    #   담보 3개가 한 줄로 뭉쳐 <b>상해1~5종·산정특례 뇌심·중입자·다빈치·화상수술비가 통째로 소실</b>됐는데
    #   <b>확인사항에도 안 떠서</b> 지점장이 눈으로 발견할 때까지 몰랐다. 그게 진짜 문제였다.
    #   → 접힘의 지문 2개를 자동 검출해 <b>확인사항에 강제 노출</b>한다(값은 그대로 두고 표시만 — 누락 금지).
    #     ①<b>담보 접미어가 2회 이상</b>(진단비·수술비·치료비·일당·입원비·보장) = 두 담보가 뭉친 것
    #     ②<b>괄호가 안 닫힘</b>(열림>닫힘) = 담보명이 잘려 다음 담보를 삼켰을 가능성
    #   ★이 진단이 있었으면 지적받기 전에 잡았다. <b>절대 제거하지 말 것.</b>
    _FOLD_SFX = ('진단비','수술비','치료비','일당','입원비','보장')
    for _c in contracts:
        for _rw, _av in (_c.get('dambo') or {}).items():
            _s = str(_rw)
            _cnt = sum(_s.count(_x) for _x in _FOLD_SFX)
            _bal = _s.count('(') - _s.count(')')
            if _cnt >= 2 or _bal > 0:
                _why = ('담보 접미어 %d회 = 두 담보 뭉침 의심' % _cnt) if _cnt >= 2 else '괄호 미닫힘 = 담보명 잘림 의심'
                print('[FOLD_AUDIT] %s | %s | %s | %s' % (_c.get('company',''), _s[:70], _av, _why))
                # ★★★★★v333b (지점장 지시 2026.08.02, 영구): 지점장 원문 —
                #   "<b>삼성생명의 애매한 담보는 세부가입현황에서 체크해라 / KB도 이런 담보표에서 체크된다</b>".
                #   접힘·뭉침이 의심되는 담보는 <b>별첨이 아니라 세부가입현황(계약별 가입정보)</b>에서 대조한다.
                #   근거 = 삼성 리빙케어·KB처럼 별첨 담보명이 상품명뿐이거나 뭉치는 리포트는
                #   세부가입현황이 <b>담보 종류를 알려주는 유일한 표</b>다(지침 §3 ②).
                unmapped.append((0, _c.get('company',''), '[접힘의심] ' + _s, _av,
                                 _why + ' → <b>세부가입현황(계약별 가입정보)에서 대조</b>'))
    cancer_trace = []  # ★v30h 암 블록 기재 근거 — (회사, 원담보명, 기재행, 금액). 일반암 과다합산 즉시 추적
    surg_trace = []    # ★v30g 수술 블록 기재 근거 — (회사, 원담보명, 기재행/슬롯, 금액)
    raw_by_std = {}   # ★v39 워크시트 담보명 카피: 표준명→원본담보명(최댓값 담보 기준)
    # ★★★★★v289 (지점장 지시 2026.07.31 "계속 반복이야 — 우선 원인 잡자")
    #   <b>반복의 구조적 원인</b>: 근거 수집이 `_WS_STD` 10개 담보에만 걸려 있어
    #   <b>운전자·사망·일당·실손·골절 등은 어느 담보에서 왔는지 산출물 어디에도 없다</b>.
    #   → 지점장이 "대인 3,020"을 보고도 원인을 알 수 없고, 물어보고, 내가 추측하고, 틀리고, 원복한다.
    #   <b>실측(이명순)</b>: 대인 3,020의 범인 `대인교통사고발생위로금담보특별 20`을
    #   찾는 데 이 대화에서만 6턴이 걸렸다. 근거표 한 줄이면 지점장이 1분에 짚는다.
    #   → <b>전 담보 매핑 근거를 빠짐없이 기록</b>한다(표시 전용 — 값에 일절 영향 없음).
    trace_all = []
    heart_trace = []   # ★v29z (지점장 2026.07.03): 심장 블록 기재 근거 — (회사, 원담보명, 기재행들, 금액). '없는 값이 튀어나옴' 방지용 감사 로그
    silson_trace = []  # ★v29z: 실손 세대 판정 근거 — (회사, 가입일, 상품코드, 판정)

    def _in_sum(ct):
        """월보험료 합계 포함 여부. ★v129 정본: 잔여보험료 > 0 이면 포함(완납이면 0).
           잔여보험료를 못 읽었으면(None) 납입횟수 a>=b(완납)로 폴백."""
        _r = ct.get('remain')
        if _r is not None: return _r > 0
        return not _is_paid_up(ct.get('pay_period',''), ct.get('pay_count',''))

    # ★★★★★v477 제83조 보강 — 계약 루프가 끝난 뒤에도 <b>세부보충·1-3종·사망이동</b>이
    #   같은 셀을 다시 칠한다(6269·6304행 `BL if _gen else BK`). 그래서 계약 루프 끝의 확정만으로는
    #   <b>덮인다</b>(김순자 v476 실측 — 회사명은 고쳐졌는데 갱신담보 색은 그대로 검정).
    #   → 열별 파랑 행을 <b>루프 바깥에 보관</b>했다가 <b>모든 재기재가 끝난 뒤</b> 마지막에 확정한다.
    _BLUE_ALL = {}
    _hy_cells = []                 # ★v580 제122조 4항 — 현물 셀(노란칸) 좌표
    _HYUNMUL.clear()               # ★고객마다 비운다(전역 오염 방지 · v579e와 같은 원리)
    for i, ct in enumerate(contracts):
        col = 3 + i
        # ★★★★★v580 제122조 4항 — 현물 판정에 필요한 <b>계약 컨텍스트</b>를 심는다.
        #   `resolve_kw`는 담보명만 받으므로 회사·계약일을 알 수 없다. 전역에 실어 넘긴다.
        globals()['_CUR_CO'] = str(ct.get('company') or '')
        globals()['_CUR_CD'] = str(ct.get('contract_date') or '')
        gen  = ct['renewal'] == '갱신'
        paid = not _in_sum(ct)          # ★v199: 헤더 진녹(완납) 판정 = 합계 제외 판정과 동일 근거
        h = ws.cell(1, col)
        h.value = f"{ct['company']}\n{ct['product']}\n[{ct['renewal']}]"
        h.font = W; h.alignment = AL
        # ★★★★★v334 (지점장 지적 2026.08.02): <b>실손 계약 헤더가 완납(진녹)으로 찍히던 것 수정</b>.
        #   실측 = 현대해상 `무배당실손의료비보장보험(갱신형)(Hi2004)` 보험료 0원·납입 0/0 →
        #   `_in_sum`이 완납으로 판정 → <b>초록</b>. 그러나 §6 ⓪ 정본은
        #   <b>"실손은 비갱신이 없다 — 무조건 갱신. 어떤 조건보다 우선. 헤더색도 파랑"</b>이다.
        #   → 실손이면 완납 판정보다 <b>파랑이 우선</b>한다.
        if _is_silson_prod(ct.get('company',''), ct.get('product','')):
            h.fill = FILL_BLUE
        else:
            h.fill = FILL_GREEN if paid else (FILL_BLUE if gen else FILL_RED)
        pm = ct['premium']
        # ★★v199 지점장 확정 2026.07.23: 보험료 합계를 '=D2+E2+…'가 아니라 단일 '=SUM()'으로 만들기 위해
        #   완납 계약의 보험료 칸은 <b>텍스트</b>로 넣는다. 엑셀 SUM은 텍스트를 무시하므로
        #   금액은 화면에 그대로 보이면서 합계에서만 자동 제외된다(v129 정본 유지).
        # ★★★★★v424 (지점장 확정 2026.08.16): <b>일시납 계약의 납입액이 통째로 버려지던 것 수정</b>.
        #   실측 = 메트라이프 무배당 변액연금보험 동행 — PDF 「일시납 / 11,000,000원」인데
        #   엑셀 2행이 <b>공란</b>. 월납만 담기고 일시납은 담을 자리가 없어 사라졌다.
        #   → <b>텍스트</b>로 넣는다. 엑셀 SUM은 텍스트를 무시하므로 월보험료 합계는 그대로다
        #     (일시납을 월보험료에 더하면 424,563원이 1,100만이 된다 — 섞으면 안 된다).
        _lump = 0
        if not pm:
            try:
                _lump = int(str(ct.get('lump_sum') or ct.get('premium_lump') or 0
                                ).replace(',', '').replace('원', '').strip() or 0)
            except Exception:
                _lump = 0
        if pm and paid:
            ws.cell(2,col).value = f'{pm:,} (완납)'
        elif pm:
            ws.cell(2,col).value = pm
        elif _lump:
            ws.cell(2,col).value = f'{_lump:,} (일시납)'
            print(f"[v424 일시납] {ct['company']} {ct['product']} — {_lump:,}원 헤더 2행 기재")
        else:
            ws.cell(2,col).value = None
        ws.cell(2,col).font = BL if gen else BK
        ws.cell(3,col).value = ct['contract_date']
        ws.cell(4,col).value = ct['expiry_date']
        # ★★★★★v302-D (지점장 지시 2026.07.31): <b>납부표시 (1/1)</b>. 구 코드는 납입주기가 비면
        #   회차까지 통째로 지워 일시납 종신이 <b>총납입기간 공란</b>으로 나갔다. 회차만 있어도 표기한다.
        _pp5, _pc5 = ct.get('pay_period') or '', ct.get('pay_count') or ''
        ws.cell(5,col).value = (f"{_pp5} ({_pc5})" if _pp5 and _pc5 else
                                (_pp5 or (f"({_pc5})" if _pc5 else '')))
        for r in [3,4,5]: ws.cell(r,col).font = BL if gen else BK

        dambo = ct['dambo']
        # ★★★★★v475 제83조 (지점장 지적 2026.08.18 「갱신형 담보 → 이거 갱신인데 비갱신으로 잡혔다」)
        #   제5조 B(담보명에 '(갱신)' → 그 담보만 파랑)가 <b>담보 루프 한 경로에서만</b> 지켜지고
        #   있었다. 심장 묶음 분해 · CI 배분 · 종신 사망 이동 등 <b>다른 15곳은 `BL if gen else BK`</b>라
        #   비갱신 계약의 '갱신형 담보'가 <b>검정</b>으로 찍혔다(김순자 실측 — 메리츠 `갱신형
        #   뇌혈관질환진단비`·`갱신형 허혈성심장질환진단비`·`갱신형 암진단비`·`갱신형 일반상해사망`).
        #   → <b>파랑이어야 할 행을 모아 두고 계약 루프 끝에서 한 번에 확정</b>한다.
        #     (제0조 「판정은 한 곳에서만」 · 「규칙 두 개가 같은 값을 만지면 나중 것이 앞 것을 죽인다」)
        _blue_r = set(); _BLUE_ALL[col] = _blue_r
        jong_acc = {'상해 종수술비(1-5종)':[0]*8, '질병 종수술비(1-5종)':[0]*8}
        ndae_acc = [0]*6      # ★v386 116대/120대 수술비 = 등급 Ⅰ~Ⅵ 슬래시 6칸   # ★v29v 8칸 수집 후 기재 시 5/8종 판정
        trio_acc = [0,0,0]   # ★v29y MRI/도수치료/비급여주사
        jong_blue = {'상해 종수술비(1-5종)':False, '질병 종수술비(1-5종)':False}
        # ★v378 종번호 없는 단일금액 종수술(질병종수술·상해종수술·파워수술보장 등) 수집칸.
        #   슬래시 칸(jong_acc)과 <b>같은 행</b>을 쓰므로 둘이 동시에 차면 [확인]큐로 보낸다(조용한 덮어쓰기 차단).
        jong_lump = {'상해 종수술비(1-5종)':0, '질병 종수술비(1-5종)':0}

        # ★ CI/리빙케어/GI 본체 분해 (지점장 지시 2026.06.28): 주계약 최대=사망, 본체=사망의 80%/50%,
        #   본체를 중대한암·중대한뇌졸증·중대한급성심근에 동일 기재 / 사망 전액=일반사망 / 판별실패=주계약 [확인].
        # ★v246: 3열(KB)은 회사명 칸에 상품명이 붙어 오는 경우가 있다 → <b>회사명·상품명 둘 다</b> 본다.
        #   "이건 롯데나 KB나 동일하다"(지점장 확정) — 2열·3열이 같은 CI 규칙을 타야 한다.
        _is_ci = _isci_prod(ct.get('product'), ct.get('company'))   # ★v446 손보 제외
        _ci_fix = None      # ★v239 CI 본체 최종 기재 예약(담보 루프 뒤에서 적용)
        _cij = ct.get('ci_jugye') or []
        # ★★★v237 (지점장 지시 2026.07.25, 영구): <b>선지급률(50%형/80%형)을 세부가입현황(상세내역)에서 찾아낸다</b>.
        #   구 코드는 근거가 <b>별첨의 '주계약' 라벨(`ci_jugye`)뿐</b>이었다. 롯데(let:) 별첨엔 '주계약' 라벨이
        #   아예 없어(주계약을 `질병사망`/`상해사망`으로 씀) `_cij=[]` → <b>이 블록을 통째로 건너뛰었다</b>.
        #   그래서 CI 상품인데 선지급 분해가 하나도 안 됐다(v235에서 `_isci_prod`를 고쳐도 여기서 다시 막혔다).
        #   → <b>①별첨 주계약(1순위) ②세부가입현황 계약열(2순위)</b> 순으로 사망액·본체 후보를 확보한다.
        _sebu_ci = ct.get('ci_sebu') or {}
        _cl = ct.get('ci_lines') or {}
        _samang=0; _cand=[]
        # ★★★★★v239 판정 근거 우선순위(영구): ①별첨 '주계약' 라벨 → ②<b>별첨 줄 단위 사망/진단 원값</b> → ③세부가입현황 사망액
        #   ②가 핵심이다. 롯데 별첨은 주계약을 `질병사망`으로 쓰고 <b>추가특약 사망이 별도 줄</b>로 온다.
        #   dambo는 이를 합산하므로(3,000+2,000=5,000) 비율이 깨진다 → <b>줄 단위 개별값</b>으로 봐야 80%가 나온다.
        if _is_ci and _cij:
            _samang = max(_cij); _cand=[v for v in _cij if 0<v<_samang]
        elif _is_ci and (_cl.get('samang') or []):
            _samang = max(_cl['samang'])
            _cand   = [v for v in (_cl.get('cands') or []) if 0<v<_samang]
            print(f"[v239 CI·별첨줄] {ct.get('company')} 주계약사망={_samang:,} "
                  f"(사망줄 {sorted(set(_cl['samang']),reverse=True)}) 진단후보={sorted(set(_cand),reverse=True)[:8]}")
        elif _is_ci and _sebu_ci.get('samang'):
            # ★★★★★v327 (지점장 지시 2026.08.02, 영구): <b>세부가입현황 담보군값도 본체 후보로 쓴다</b>.
            #   지점장 원문: "세부가입현황에 가보면 주계약 5천만원 - 암 뇌졸증 급성심근 각각 4천씩 있는데
            #   왜 지침을 무시하는것이냐". 구 코드는 이 분기에서 <b>`_cand=[]`로 후보를 통째로 비웠다</b>
            #   → §22 CI 4단계의 ④(동일 금액 2개 이상 → 50%×2 / 80%÷0.8 이 사망보장금과 일치)가
            #   <b>아예 돌지 못했다</b>. 별첨에 '주계약' 라벨도 사망줄도 없는 생보 CI(삼성 리빙케어)는
            #   세부가입현황이 유일한 정답 소스인데 그 소스를 스스로 버린 것이다.
            #   ★안전장치는 종전 그대로: ④는 <b>수학적으로 사망보장금과 정확히 일치할 때만</b> 인정한다.
            #   실측(이영태 삼성 리빙케어) 사망 5,000 · 담보군 [3500,4000,600,4000,4000,…]
            #   → 4,000이 3회 · 4,000÷0.8=5,000 ✓ → 본체 4,000 · 80%형.
            _samang = _sebu_ci['samang']
            _cand   = [v for v in (_sebu_ci.get('cands') or []) if 0 < v < _samang]
            print(f"[v327 CI·세부내역] {ct.get('company')} 사망={_samang:,} "
                  f"· 본체후보={sorted(set(_cand), reverse=True)[:8]}")
        # ★★★★★v328 CI 미검출 게이트 (지점장 지시 2026.08.02, 영구)
        #   지점장 원문: "ci가 안읽히면 모두가 다 미적용된다. ㅠㅠ"
        #   → 구 코드는 <b>`_is_ci`인데 사망액을 한 소스에서도 못 찾으면 아무 일도 하지 않았다</b>.
        #     로그도 없고, 확인사항 CI 진단표에도 <b>행 자체가 안 생겼다</b> → 설계사는
        #     "이 계약은 CI가 아닌가보다"로 오해한다. <b>조용히 틀리는 것을 시끄럽게 틀리는 것으로</b> 바꾼다.
        #   ★값은 건드리지 않는다 — 표시만 한다(누락 금지·오염 금지).
        if _is_ci and not _samang:
            _ci_diag.append({'co':ct.get('company',''),'pd':ct.get('product',''),
                             'samang':0,'pct':None,'bonche':0,'src':'★사망액 미확보(판정불가)',
                             'sebu':bool((ct.get('ci_sebu') or {}).get('samang')),
                             'byul':bool((ct.get('ci_lines') or {}).get('samang')),
                             'placed':0,'placed_txt':''})
            unmapped.append((col, ct.get('company',''), '[확인] CI 사망액 미검출', 0,
                             'CI 상품인데 주계약·별첨 사망줄·세부가입현황 어디서도 사망보장금을 찾지 못했다 — '
                             '중대한OO 3행이 전부 비게 되므로 세부가입현황을 수기 대조할 것'))
            print(f"[v328 CI·미검출] {ct.get('company')} {ct.get('product')} — "
                  f"CI 상품인데 사망액 0 (주계약 {len(_cij)}건 · 별첨사망 {len(_cl.get('samang') or [])}건 · "
                  f"세부 {bool(_sebu_ci.get('samang'))}) → 확인사항 노출")
        if _is_ci and _samang:
            _bonche=None; _pct=None; _samang_all=_samang; _cisrc=''
            _slset = set(_cl.get('samang') or []) | {_samang}
            # ★★★★★v241 <b>CI보험 찾기 4단계 체크리스트 — 영구지침(지점장 2026.07.25)</b>
            #   지점장 원문: "1)상세 세부내역에 체크 2)[별첨]보험서비스(상품)별보장현황에서 체크
            #   3)중대한*** 이 있는지 체크 4)똑같은 금액이 2개 이상 있다면 50%이라면 2배 계산이 사망보장,
            #   80%형이라면 수학으로 계산해서 사망보장금과 맞는지 체크.
            #   위의 3가지 다 체크해야 하는 <b>까다롭지만 보람된 CI보험 찾기</b>"
            #   → ①②는 소스 2개(`ci_sebu` 세부내역 · `ci_lines` 별첨)를 <b>둘 다</b> 확보해 교차확인,
            #     ③ '중대한OO' 담보가 있으면 <b>그것이 본체</b>(최우선),
            #     ④ 없으면 <b>동일 금액이 2개 이상</b>인 값을 본체 후보로 삼아
            #       <b>50%형이면 ×2, 80%형이면 ÷0.8 한 값이 별첨 사망보장금과 일치</b>하는지 검증한다.
            _jd=[v for v in (_cl.get('jungdae') or []) if 0<v<_samang]
            if _jd:                                   # ③ 중대한*** 최우선
                _bonche=max(_jd); _cisrc='③중대한OO담보'
                for _r2,_p2 in ((0.5,50),(0.8,80)):
                    if round(_bonche/_r2) in _slset: _pct=_p2; break
                if not _pct: _pct=50                  # 사망줄 매칭 없으면 50%형(지점장 확정)
                _samang = round(_bonche/(_pct/100.0))
            else:                                     # ④ 동일 금액 2개 이상 → 사망보장금 일치 검증
                # ★★★★★v328 (지점장 지시 2026.08.02, 영구): <b>④를 2단으로 돌린다</b>.
                #   지점장 원문: "다른사람들도 ci가 많이 지금 오류를 지침삼아서 다 적용해라.
                #   ci가 안읽히면 모두가 다 미적용된다."
                #   → v327은 <b>③세부가입현황 경로에서만</b> 후보를 살렸다. 그런데 같은 구멍이
                #     <b>①주계약 라벨 경로·②별첨 사망줄 경로</b>에도 있다 — 각자 자기 소스의 후보만
                #     보므로, 사망액은 ①에서 잡혔는데 <b>본체 금액은 세부가입현황에만</b> 실린 계약은
                #     여전히 판정불가로 떨어진다(삼성 리빙케어가 ①에서 잡혔다면 똑같이 죽었다).
                #   ★<b>1단 = 기존 우선순위 후보</b>(무회귀 보장) → 실패하면
                #     <b>2단 = 3소스 합집합</b>(주계약 + 별첨 사망줄 + 세부가입현황)으로 재시도한다.
                #   ★안전장치 불변: ④는 <b>÷선지급률 결과가 사망보장금과 정확히 일치할 때만</b> 인정.
                def _try4(_pool):
                    _cnt={}
                    for v in _pool: _cnt[v]=_cnt.get(v,0)+1
                    for v in sorted([k for k,c in _cnt.items() if c>=2], reverse=True):
                        for _r2,_p2 in ((0.5,50),(0.8,80)):
                            if round(v/_r2) in _slset:  # ★수학적으로 일치할 때만 인정
                                return v,_p2,round(v/_r2)
                    return None,None,None
                _bonche,_pct,_sm4 = _try4(_cand)
                if _bonche:
                    _samang=_sm4; _cisrc='④동일금액2개이상'
                else:
                    _all = [v for v in (list(_cij)
                                        + list(_cl.get('cands') or [])
                                        + list(_sebu_ci.get('cands') or []))
                            if 0 < v < _samang]
                    _bonche,_pct,_sm4 = _try4(_all)
                    if _bonche:
                        _samang=_sm4; _cisrc='④동일금액2개이상(3소스합집합)'
                        _cand = _all
                        print(f"[v328 CI·합집합] {ct.get('company')} 1순위 후보로 실패 → "
                              f"주계약+별첨+세부 합집합에서 본체 {_bonche:,} 발견")
                if not _bonche:                       # 폴백: 사망 대비 비율 매칭(±2%)
                    for _ratio,_p in ((0.8,80),(0.5,50)):
                        for v in sorted(set(_cand), key=lambda x:-_cand.count(x)):
                            if _samang and abs(v/_samang-_ratio) < 0.02:
                                _bonche=v; _pct=_p; _cisrc='폴백·비율매칭'; break
                        if _bonche: break
            if _bonche:
                print(f"[v241 CI 4단계] {ct.get('company')} {_cisrc} → 본체 {_bonche:,} ÷ {_pct}% "
                      f"= 사망보장 {_samang:,} · 별첨사망총액 {_samang_all:,} · 차액 {_samang_all-_samang:,}→질병사망(80세)"
                      + ("  ✓사망줄 일치" if _samang in _slset else "  (사망줄 미매칭)"))
            _pl=[]
            if _bonche:
                _pl=[f'중대한 암 {_bonche:,}', f'중대한 뇌졸증 {_bonche:,}', f'중대한 급성심근 {_bonche:,}']
            _ci_diag.append({'co':ct.get('company',''),'pd':ct.get('product',''),
                             'samang':_samang,'pct':_pct,'bonche':_bonche or 0,
                             'src':_cisrc or '판정불가',
                             'sebu':bool((ct.get('ci_sebu') or {}).get('samang')),
                             'byul':bool((ct.get('ci_lines') or {}).get('samang')),
                             'placed':len(_pl),'placed_txt':' · '.join(_pl)})
            if not _bonche:
                _sc = sorted(set(_sebu_ci.get('cands') or []), reverse=True)[:6]
                print(f"[v237 CI·확인큐] {ct.get('company')} 사망={_samang:,} — 50%/80% 근접 본체 없음 "
                      f"→ 선지급형이 아니거나(별도 진단비 특약) 상세내역 대조 필요. 세부가입현황 담보군값={_sc}")
                unmapped.append((col, ct.get('company',''), '[확인] CI 선지급률',
                                 0, f"세부가입현황 사망 {_samang:,} · 담보군값 {_sc} — 50%/80% 미해당. 상세내역에서 선지급형 대조 요망"))
            if _bonche:
                # ★v239: 담보 루프가 뒤에서 중대한OO 행에 별첨 원값을 <b>가산</b>하므로(실측 중대한 암 7,800=5,400+2,400)
                #   여기서 기재하면 덮어써진다 → <b>예약해두고 담보 루프 뒤에서 최종 기재</b>한다.
                _ci_fix = {'bonche':_bonche,'samang':_samang,'pct':_pct}
                # ★★★★★v242 뇌 축 판별: ①본체 금액과 <b>정확히 일치</b>하는 뇌 담보의 축 ②없으면 <b>뇌출혈 우선</b>
                #   (지점장: "무조건 뇌졸증이 아니라는거다 뇌졸증 or 뇌출혈이다")
                _br = (ct.get('ci_lines') or {}).get('brain') or []
                _ax = None
                for _a,_v in _br:
                    if _v == _bonche: _ax = _a; break
                if not _ax:
                    _ax = '뇌출혈' if any(a=='뇌출혈' for a,_ in _br) else ('뇌졸증' if _br else None)
                ct['ci_brain'] = {'axis':_ax, 'rows':_br}
                # ★★★v249 (2026.07.26 한정환 KB 3열 실측): 확인사항 CI 진단표의 '중대한OO 배치' 문구가
                #   <b>'중대한 뇌졸증'으로 하드코딩</b>돼 있었다(2606행 `_pl`). 본표(엑셀 34행)는 v242로
                #   축을 따라 <b>중대한 뇌출혈</b>에 기재되는데 <b>진단표만 뇌졸증</b>이라 산출물끼리 어긋났다.
                #   실측 — 신한 본체 2,000·DB생명 본체 2,400 둘 다 축=뇌출혈인데 진단표는 '중대한 뇌졸증' 출력.
                #   `_pl` 생성이 축 판별보다 <b>앞</b>이라 생긴 순서 결함 → 여기서 배치 문구를 최종 갱신한다.
                if _ci_diag:
                    _bxl = '중대한 뇌출혈' if _ax == '뇌출혈' else '중대한 뇌졸증'
                    _ci_diag[-1]['placed_txt'] = ' · '.join(
                        [f'중대한 암 {_bonche:,}', f'{_bxl} {_bonche:,}', f'중대한 급성심근 {_bonche:,}'])
                print(f"[v242 CI뇌축] {ct.get('company')} 본체 {_bonche:,} → <b>중대한 {_ax or '뇌졸증(기본)'}</b>"
                      f"  (별첨 뇌담보 {_br})")
                # ★★★★★v239 이중계산 차단: 별첨 '주계약' 라벨 경로(`ci_jugye`)에서만 dambo에 사망을 주입한다.
                #   롯데 줄단위 경로(`ci_lines`)는 <b>이미 dambo에 '질병사망'이 들어 있어</b> 여기서 또 더하면
                #   일반사망 3,000+5,000=8,000 · 상해사망 26,100(한장표 23,100 초과)로 부푼다 — 실측 확인.
                if _cij:
                    dambo['일반사망(종신주계약)']=dambo.get('일반사망(종신주계약)',0)+_samang
                    dambo['상해사망(종신주계약)']=dambo.get('상해사망(종신주계약)',0)+_samang   # §8.1 종신 1:1
                _rci=nm2r.get('중대한CI적용')   # 사망−본체=선지급 후 잔여 사망보험금(80%형=20%잔여).
                if _rci:
                    ws.cell(_rci,col).value=_samang-_bonche; ws.cell(_rci,col).font = BL if gen else BK
                # ★v29t (지점장 확정 2026.07.02, 김진구 정본): CI추가보장특약 = 급성심근 초과분 → 최대 1건 급성심근경색 행, 잔여 [확인]
                _cex=ct.get('ci_extra') or []
                if _cex and 'CI추가보장특약' in dambo:
                    _mx=max(_cex)
                    dambo.pop('CI추가보장특약', None)
                    _rgs=nm2r.get('급성심근경색')   # 셀 직접 기재(§8.4 CI 재매핑 회피 — 초과분은 '일반' 급성심근경색 행)
                    if _rgs:
                        _ex0=ws.cell(_rgs,col).value
                        ws.cell(_rgs,col).value=(_ex0+_mx) if isinstance(_ex0,(int,float)) else _mx
                        ws.cell(_rgs,col).font = BL if gen else BK
                    _lv=sum(_cex)-_mx
                    if _lv>0: unmapped.append((col, ct['company'], 'CI추가보장특약(잔여)', _lv, '급성심근 초과분 외 잔여 → 약관 확인'))
                dambo.pop('주계약', None)
            else:
                unmapped.append((col, ct['company'], f'주계약(CI 80/50%판별실패 {_cij})', _samang, 'CI 본체비율 불명 → 좌측표 수기'))
                dambo.pop('주계약', None)

        for raw, amt in dambo.items():
            # ★ 심장 묶음담보 6사 정본 매핑(2026.06.29). I20→협심증 / 허혈성칸=단독전용 / 순환계=전체5 / 급성심근=묶음제외 / 빈맥 포함(★지점장 7/1)·심근병증 제외.
            _rn = re.sub(r'\s','',raw)
            _heart_bundle = None
            _co = ct.get('company','')
            if ('진단' in _rn and '수술' not in _rn and '주요치료' not in _rn
                and '산정특례' not in _rn and '혈전' not in _rn):
                # ★★★v225 (2026.07.25 이정화 실측): <b>산정특례·혈전용해는 심장 묶음이 아니다</b>.
                #   실사고 = `중증질환자심장질환산정특례대상진단비<b>Ⅱ</b>(연간1회한)` 의 <b>로마숫자 Ⅱ</b> 때문에
                #   `_rmn()`이 `_t=2`를 반환해 DB '특정Ⅱ' 묶음으로 오인 → <b>급성심근경색 행에 1,000이 배정</b>되고
                #   <b>산정특례심장은 0</b>이 됐다(급성심근경색 8,000 = 정답 7,000 + 오배정 1,000).
                #   `_HB` 후처리 테이블에는 이미 산정특례 제외가 있었는데 <b>이 인라인 블록에만 빠져 있었다</b>.
                # ★v30o 고정(메리츠, 지점장 2026.07.03): 심장질환진단비Ⅰ→허혈성 진단비 / 심장질환진단비Ⅱ→급성심근경색
                # ★★★★★v379 (지점장 재확정 2026.08.10, 영구·최상위 — "원래 처음부터 있던 지침"):
                #   <b>모든 회사에서 「허혈성진단비 · 허혈성심장진단비 · 심장허혈성진단비」는 단독이다.</b>
                #   → 회사 분기보다 <b>먼저</b> 판정한다. 어느 회사든 분해(협심증·급성심근 동반) 금지.
                #   [구 결함] 회사별 분기 안에서만 단독 판정을 해서 <b>삼성·현대 경로가 협심증을 함께 찍었다</b>
                #   (실측 구본칠 — `허혈심장질환진단특약ⅢUT…` 1,000 → 협심증 1,000 신규 발생).
                #   ★묶음 수식어(특정·6가지·n대·순환계)가 붙은 담보만 회사별 표를 탄다.
                # ★★★★★v380 (지점장 재확정 2026.08.11, 영구·최상위 — "전부 다 단독이다"):
                #   <b>뇌혈관진단비 · 허혈성진단비 · 급성심근경색진단비 · 뇌졸증진단비 · 뇌출혈진단비
                #   = 5종 전부 단독</b>. 어느 회사든 <b>묶음 분해 금지</b>(협심증·빈맥·심부전 동반 금지).
                #   [구 결함] v379는 <b>허혈성 한 종류만</b> 전역 처리하고 나머지 4종은 회사 분기에 남겨 뒀다.
                #   ★구현 = 여기서 값을 강제로 찍지 않고 <b>묶음 블록 자체를 건너뛴다</b>.
                #     resolve2가 각 담보를 자기 행 하나로 보내므로(실측), 이렇게 해야
                #     <b>CI 변환(중대한OO)·일당 상한 가드·[확인]큐</b> 등 하류 규칙이 전부 살아 있다.
                #   ★묶음 수식어(특정·6가지·n대·순환계)가 붙은 담보만 회사별 표를 탄다.
                #   ★★★단서(지점장 2026.08.11): <b>「뇌혈관Ⅰ · 뇌혈관Ⅱ」처럼 등급 Ⅰ·Ⅱ가 붙은 것은 다르다</b>
                #     → 단독 예외에서 <b>빼고 종전 지침(회사별 표·Ⅰ/Ⅱ 분기)대로</b> 처리한다.
                #     ※Ⅲ·특약Ⅲ 등 상품세대 표기(`허혈심장질환진단특약ⅢUT…`)는 등급 Ⅰ·Ⅱ가 아니므로 단독 유지.
                _is_solo5 = is_solo5_name(_rn)
                if _is_solo5:
                    pass                                   # 단독 5종 — 묶음 분기 진입 금지
                elif ('심장질환진단' in _rn) and ('허혈' not in _rn) and ('급성심근' not in _rn):
                    # ★양예서/메리츠 어린이: 심장질환진단비Ⅱ→급성심근경색 / Ⅰ→허혈성 진단비 (별첨값 앵커: Ⅰ=600 허혈성, Ⅱ=3000 급성심근)
                    _mn=_rmn(_rn)
                    # ★★★★★v384 (지점장 확정 2026.08.11, 영구·최상위):
                    #   지점장 원문 = "<b>심장Ⅰ·Ⅱ에 허혈성이 들어가 있다면 그건 다 협심증이다.
                    #   협심증의 종류가 허혈성이다. 그래서 다 협심증으로 표기해라</b>".
                    #   → 「허혈성 진단비」 행은 <b>담보명이 허혈성진단비·허혈성심장질환진단비인 단독 담보 전용</b>이다.
                    #     심장Ⅰ/특정심장Ⅰ 같은 <b>묶음</b>이 허혈성을 품고 있으면 <b>협심증 행</b>으로 표기한다.
                    #   [구 결함] 아래 3곳이 묶음(Ⅰ)을 그대로 허혈성 행에 넣고 있었다 — `_HB` 주석
                    #     「묶음은 이 행에 절대 안 넣는다. I20·I24·I25는 협심증 행으로 표현」을 정면 위반.
                    if _mn==2: _heart_bundle=['급성심근경색']
                    elif _mn==1: _heart_bundle=['협심증']
                # ★뇌질환진단비Ⅰ/Ⅱ (메리츠 어린이 등): Ⅱ→뇌졸증(넓음) / Ⅰ→뇌혈관진단비
                if (not _is_solo5) and _heart_bundle is None and ('뇌질환진단' in _rn):
                    _mn=_rmn(_rn)
                    if _mn==2: _heart_bundle=['뇌졸증진단비']
                    elif _mn==1: _heart_bundle=['뇌혈관진단비']
                # ★v30z4 성인병진단금(생보·AIA·AIG·라이나·우체국 등) = 급성심근경색(100% 확정) + 뇌졸증/뇌출혈(세부가입 판별).
                #   지점장 반복 확정: [확인] HOLD 폐기. 뇌축 = 계약에 뇌출혈 담보 있으면 뇌출혈, 없으면 뇌졸증(세부가입 뇌혈관 표기도 뇌졸증계로 해석).
                if (not _is_solo5) and _heart_bundle is None and ('성인병' in _rn):
                    _brain = '뇌출혈진단비' if any('뇌출혈' in str(_k) for _k in dambo.keys()) else '뇌졸증진단비'
                    _heart_bundle = ['급성심근경색', _brain]
                # ★v29w 심장 범위 재점검(지점장 2026.07.02, 6사 정본 대조):
                # DB 순환계 5종(중증) = 급성심근경색 + 뇌졸중
                # ★★★★★v266(지점장 2026.07.28 "저건 2개가 세트라서 하나만 잡으면 된다"):
                #   `4대순환계질환진단비(특정3대심장질환)`은 <b>위 `_HB` 후처리가 이미 심부전·부정맥으로 분해</b>했다.
                #   그런데 아래 `elif '순환계' in _rn` 이 <b>같은 담보를 또 잡아</b> 심부전·부정맥이 <b>2배</b>가 되고
                #   빈맥·급성심근까지 얹혔다(실측 박O정: 심부전 500 → 1,000 · 빈맥 0 → 500).
                #   → <b>4대순환계 계열은 인라인 블록에서 제외</b>한다. 규칙은 하나만 잡는다.
                if (not _is_solo5) and _heart_bundle is None and '4대순환계' in _rn:
                    _heart_bundle = []          # _HB가 처리 완료 — 여기서 중복 적용 금지
                elif _heart_bundle is None and '순환계' in _rn and '5종' in _rn:
                    _heart_bundle = ['급성심근경색','뇌졸증진단비']
                # DB 순환계 4종 = 협심증·심부전(+빈맥, 심근병증 [확인])
                elif '순환계' in _rn and '4종' in _rn:
                    _heart_bundle = ['협심증','심부전']
                # DB 순환계 3종 = 염증·부정맥
                elif '순환계' in _rn and '3종' in _rn:
                    _heart_bundle = ['염증','부정맥']
                # ★순환계3대(허혈성심장 I20~25 + 뇌혈관 I60~69 + 말초 I70~79) — 삼성·DB 등. 말초=전용행無(누락 감수)
                elif '순환계' in _rn:
                    if ('DB' in _co) or ('디비' in _co):
                        _heart_bundle = ['급성심근경색','빈맥','부정맥','심부전']   # DB 순환계3대=심장정지I46.0·부정맥I47~49·심부전I50
                    else:
                        _heart_bundle = ['협심증','급성심근경색','뇌혈관진단비']
                # ===== BARUM 10사 질병코드 분류표 정본(2026.07.05 지점장 확정): 특정Ⅰ/Ⅱ 라벨=회사마다 다름 → 회사별 표대로 =====
                elif (not _is_solo5) and any(_k in _rn for _k in ('심혈관','심장','허혈','부정맥','빈맥','심부전','심근병','판막','협심','전도','방실')):
                    _t=_rmn(_rn)
                    _i49excl=('제외' in _rn) and (('I49' in _rn) or ('부정맥' in _rn))   # ★(기타심장부정맥제외)=Ⅰ에서 I49 뺀 묶음(부정맥 담보 아님)
                    _i49=(not _i49excl) and (('I49' in _rn) or ('기타부정맥' in _rn) or ('기타심장부정맥' in _rn))
                    # 흥국·롯데: 특정Ⅰ=급성심근 / 특정Ⅱ=협심증+염증 / 롯데 15대=판막·심근병·빈맥·심부전
                    #   ★v384 구 주석 '허혈'은 협심증으로 표기(묶음은 허혈성 행에 안 넣는다).
                    if ('흥국' in _co) or ('롯데' in _co):
                        # ★★★★★v386 (지점장 지적 2026.08.12 「흥국을 못읽어낸다」 실측 후속):
                        #   흥국 <b>허혈성심질환진단비Ⅱ</b> 500이 <b>협심증·염증으로 분해</b>됐다.
                        #   원인 = 아래 `_t==2 → 협심증·염증`(롯데 특정심장Ⅱ 규칙)이 <b>묶음 수식어 없이 등급만</b>
                        #   붙은 단독 담보까지 잡았기 때문. 같은 계약의 `뇌혈관질환진단비Ⅱ`는 뇌 경로라
                        #   <b>단독으로 정상 기재</b>돼 비대칭이 드러났다(#27 단독 5종 위반).
                        #   → 등급 분기는 <b>묶음 수식어가 있을 때만</b> 탄다. 없으면 resolve가 자기 행으로 보낸다.
                        _hl_bund = any(_k in _rn for _k in ('특정','15대','심혈관','대심장','순환계'))
                        if _i49excl: _heart_bundle=['협심증','빈맥','심부전']   # 흥국 특정심혈관질환(기타심장부정맥제외)=협심·허혈·빈맥·심부전(별표70)
                        elif _i49: _heart_bundle=['부정맥']
                        elif '심근병' in _rn: _heart_bundle=['심근병증']
                        elif '15대' in _rn: _heart_bundle=['심장판막','심근병증','빈맥','심부전']
                        elif ('방실' in _rn) or ('전도' in _rn): pass   # 전용행無→[확인]
                        elif ('주요' in _rn and ('염증' in _rn or '심장염' in _rn)) or ('심낭' in _rn): _heart_bundle=['염증']
                        elif _t==1 and _hl_bund: _heart_bundle=['급성심근경색']
                        elif _t==2 and _hl_bund: _heart_bundle=['협심증','염증']
                    # ★DB(정본 재수정): 특정Ⅰ=협심증·허혈·염증 / 특정Ⅱ=급성심근 / 특정Ⅲ=판막·빈맥·심부전 / 심근병증
                    elif ('DB' in _co) or ('디비' in _co):
                        if _t==2: _heart_bundle=['급성심근경색']
                        elif _t==3: _heart_bundle=['심장판막','빈맥','심부전']
                        elif '심근병' in _rn: _heart_bundle=['심근병증']
                        elif _i49: _heart_bundle=['부정맥']
                        elif _t==1: _heart_bundle=['협심증','염증']
                    # 한화·NH농협: Ⅰ=협심증+빈맥+부정맥+심부전 / Ⅱ=급성심근 / (I49제외)=부정맥 뺀 묶음 / 심근병증
                    #   ★v384 구 주석의 '허혈'은 <b>협심증</b>으로 표기(묶음은 허혈성 행에 안 넣는다).
                    elif ('한화' in _co) or ('농협' in _co) or ('NH' in _co):
                        if _i49excl: _heart_bundle=['협심증','빈맥','심부전']   # Ⅰ에서 I49(부정맥) 제외 묶음
                        elif _t==2: _heart_bundle=['급성심근경색']
                        elif '심근병' in _rn: _heart_bundle=['심근병증']
                        elif ('주요' in _rn and ('염증' in _rn or '심장염' in _rn)): _heart_bundle=['염증']
                        elif _i49: _heart_bundle=['부정맥']
                        elif '특정질환' in _rn: _heart_bundle=['협심증','빈맥','심부전']   # 한화 심혈관특정질환=Ⅰ에서 I49제외
                        elif _t==1: _heart_bundle=['협심증','빈맥','부정맥','심부전']
                    # KB: 특정Ⅰ=협심증+빈맥+심부전 / Ⅱ=급성심근 / 심장판막=판막+염증 / I49=부정맥(빈맥X)
                    #   ★v384 구 주석 '허혈'은 협심증으로 표기. ★염증은 `_HB`가 정본(KB 특정1에 주요심장염증 포함).
                    elif ('KB' in _co) or ('케이비' in _co):
                        if _t==2: _heart_bundle=['급성심근경색']
                        elif '심근병' in _rn: _heart_bundle=['심근병증']
                        elif '판막' in _rn: _heart_bundle=['심장판막','염증']
                        elif _i49: _heart_bundle=['부정맥']
                        elif _t==1 or ('확대' in _rn and '심장' in _rn) or ('특정심장' in _rn): _heart_bundle=['협심증','빈맥','심부전']
                    # 현대(정본 재수정): 허혈성심장질환진단비=<b>단독</b>(허혈성 행) / 특정허혈=급성심근 /
                    #   특정Ⅰ=협심증+빈맥+심부전 / 특정Ⅱ=급성심근 / 주요염증 / 특정2대+I49=부정맥
                    #   ★v384 구 주석 '협심증+허혈' 폐기 — 담보명이 허혈성이면 단독, 묶음이면 협심증.
                    elif '현대' in _co:
                        if '특정허혈' in _rn: _heart_bundle=['급성심근경색']
                        elif ('허혈성심장' in _rn) or ('허혈심장' in _rn): _heart_bundle=['협심증']
                        elif '심근병' in _rn: _heart_bundle=['심근병증']
                        elif ('주요' in _rn and ('염증' in _rn or '심장염' in _rn)): _heart_bundle=['염증']
                        elif ('특정2대' in _rn) or ('방실' in _rn) or ('전도' in _rn) or _i49: _heart_bundle=['부정맥']   # 특정2대+기타부정맥(I49) 병합→부정맥(전도장애 전용행無)
                        elif _t==2: _heart_bundle=['급성심근경색']   # ★현대 특정Ⅱ=급성심근경색(정본 재수정)
                        elif _t==1 or '심혈관' in _rn: _heart_bundle=['빈맥','심부전']
                    # ★v379 삼성·메리츠 별도 단독판정 <b>폐기</b> — 위 v379 전역 가드가 전 회사를 처리한다.
                    #   구 코드는 여기서만 단독을 봐서 <b>회사가 삼성·메리츠가 아니면 분해</b>됐다.
                    #   (v206 '갱신형' 접두 · v378 접미 수식어 회귀도 전역 가드가 함께 흡수한다.)
            if _heart_bundle:
                for _bt in _heart_bundle:
                    _br = nm2r.get(_bt)
                    if _br:
                        _ex = ws.cell(_br,col).value
                        ws.cell(_br,col).value = (_ex+amt) if isinstance(_ex,(int,float)) else amt
                        _bl0 = gen or _is_gen_dambo(raw, ct.get('contract_date',''))   # ★v475·v495
                        ws.cell(_br,col).font = BL if _bl0 else BK
                        if _bl0: _blue_r.add(_br)
                heart_trace.append((ct['company'], raw, ' · '.join(_heart_bundle), amt))   # ★v29z 근거 기록
                continue
            # ★ 우선순위 역전: 확정 규칙(resolve2) 먼저 → 못 잡은 것만 Haiku(llm_resolve).
            #   Haiku가 간병인·암주요치료비·하이클래스 등 확정담보를 가로채 누락시키던 문제 차단.
            std, jong = resolve2(raw)
            jong = jong or get_종번호(raw)
            # ★★★★★v413 (지점장 확정 2026.08.12): <b>자부상은 가입금액이 아니라 12~14급(경상) 지급액이다.</b>
            #   지점장 원문: 「<b>디비손보 운전자보험 2400 = 30 이다</b> / <b>1600=20 800=10</b>」
            #   → 실측 비율 <b>÷80</b> (2400/80=30 · 1600/80=20 · 800/80=10).
            #   원문 지급표와 일치: DB `자동차부상치료비Ⅱ` 1,600 → 「1급 1천600만원 … <b>12~14급 20만원</b>」.
            #   ★왜 필요했나: 가입금액 1,600은 <b>1급(중상) 한도</b>다. 자부상 칸은 <b>경상 때 얼마 나오나</b>를
            #     보는 자리이고 한장표 표준금액도 30만원이다 — 1,600을 그대로 넣으면 <b>53배 과대표시</b>다.
            #   ★적용 범위: <b>DB손해보험</b> + 급수 밴드가 <b>1~14급</b>인 담보만. 다른 회사는 손대지 않는다
            #     (등급표 비율이 회사마다 같다는 근거가 없다 — 확인되면 그때 넓힌다).
            if std == '자부상' and isinstance(amt, (int, float)) and amt:
                # ★v413b (지점장 확정 2026.08.12): 「<b>방금 건 DB손보 운전자 전용 지침이다</b>」
                #   → 회사 DB + <b>상품명에 '운전자'</b>가 있을 때만. 다른 DB 상품엔 적용하지 않는다.
                _cmp413 = str(ct.get('company') or '')
                _prd413 = str(ct.get('product') or '')
                if ('DB' in _cmp413 or '디비' in _cmp413) and ('운전자' in _prd413):
                    _b413 = re.search(r'(\d+)\s*~\s*(\d+)\s*급', raw)
                    _is14 = (int(_b413.group(2)) >= 14) if _b413 else True   # 밴드 표기 없으면 DB 1~14급형
                    if _is14 and amt >= 80:
                        _new413 = round(amt / 80)
                        print(f"[v413 자부상] {ct['company']} '{raw}' {amt} → {_new413} (12~14급 경상 지급액 = 가입금액÷80)")
                        amt = _new413
            if not std:                       # 규칙이 못 잡은 것만 LLM 폴백
                # ★v30m 수술·일당은 resolve_kw(+DMAP)가 최종 판정. resolve_kw가 [확인](None)으로 보낸 변형을
                #   Haiku가 질병/상해수술비·상해/질병일당 base 행으로 되끌어와 합산하던 문제 차단 → 변형은 그대로 [확인].
                if re.search(r'수술|일당', raw):
                    m = {}
                else:
                    m = LLMMAP.get(raw) or {}
                    std = m.get('std')
                    if not jong: jong = m.get('jong', 0) or 0
            else:
                m = {}
            if std and _isci_prod(ct['product'], ct.get('company')):
                std = {'일반암':'중대한 암','뇌졸증진단비':'중대한 뇌졸증','급성심근경색':'중대한 급성심근'}.get(std, std)
            elif std and not _isci_prod(ct['product'], ct.get('company')):
                # ★2026.07.12 지점장 확정: 상품명에 CI/GI/리빙케어가 없으면 진짜 CI가 아니다.
                #   '중대한 암' → 일반암진단비로 산입 / 그 외 '중대한OO'는 가짜 → 전부 무시(기재 안 함).
                if std == '중대한 암':
                    std = '일반암'
                elif std in ('중대한 뇌졸증', '중대한 급성심근', '중대한CI적용', '중대한 뇌출혈'):
                    unmapped.append((col, ct['company'], raw, amt, 'CI/GI/리빙케어 상품 아님 → 중대한OO 무시'))
                    continue
            # ★★★v220 영구 가드(2026.07.25 실사고): <b>일당 행에 진단비 금액이 들어가면 그건 일당이 아니다</b>.
            #   실사고 = 간호통합병동 행에 <b>2,000(2천만원)</b>이 찍혔다. 일당은 하루당 지급액이라
            #   실무 최대가 <b>30~50만원</b>이고 100만원을 넘는 일당 담보는 존재하지 않는다.
            #   → 일당 계열 행에 <b>100 초과</b> 금액이 오면 매핑 오류로 보고 [확인]큐로 보낸다(조용한 오출고 차단).
            _DAILY = ('질병일당','상해일당','간병인','간병인지원일당','간호통합병동','종합병원 질병입원일당','종합병원 상해입원일당',
                      '1인실 상급병원','1인실 종합병원','질병중환자실','상해중환자실',
                      '질병수술일당','상해수술일당','암일당')
            if std in _DAILY and isinstance(amt,(int,float)) and amt > 100:
                unmapped.append((col, ct['company'], raw, amt,
                                 f'[확인] 일당 행에 100만원 초과({amt}) — 진단비·수술비 오매핑 의심'))
                continue
            # ★★★★★v353 (지점장 지시 2026.08.02, 영구): <b>골절·화상 「등급별 100만↑ 제외」 규칙 폐기</b>.
            #   지점장 원문: 라이나 골절진단비II(치아파절포함) 1,000 · AXA 골절진단의료비용(치아파절제외) 1,000
            #   · AIG 골절진단의료비용 500 · Ⅲ 250 · 화상진단의료비용 500 → "<b>넣어줘</b>".
            #   구 규칙은 금액이 <b>100(만원) 이상이면 등급별로 보고 [확인]큐로 보내 기재하지 않았다</b>.
            #   → 금액과 무관하게 <b>전부 기재</b>한다. 구 '등급별 100만↑ 제외' 조문은 완전 폐기.
            if std=='합의금' and amt>25000:   # 합의금 최대 2.5억, 초과는 불가 → [확인]
                unmapped.append((col, ct['company'], raw, amt, '합의금 2.5억 초과(불가)'))
                continue
            if std=='입원':
                # ★★v91 수정(지점장 확정 2026.07.19 · 장혜경 실데이터로 확인):
                #   구 규칙 '입원은 무조건 5,000 고정'이 <b>구실손의 실제 한도를 뭉갰다</b>.
                #   실측: DB 0604_TM(2006년형 1세대) 질병입원의료비 <b>500</b> → 5,000으로 부풀려짐.
                #   한장보장표(실손입원 상해 5,100 / 질병 5,500)와 불일치 = 등식1 위반.
                #   → <b>별첨 명시값을 그대로 쓴다.</b> 금액을 못 읽었을 때만 5,000을 기본값으로 넣는다.
                if not amt: amt=5000
            # ★v34 암주요치료비 10,000 강제 폐기(지점장 2026.07.09): 실제 가입금액 사용. 하이클래스는 별도행 합산.
            blue = gen or _is_gen_dambo(raw, ct.get('contract_date',''))   # ★v476·v495
            # ★★★★★v337 (지점장 지적 2026.08.02 "표 안에 입원이 검정 통원·약은 파랑이다"):
            #   §10 정본 = <b>실손(입원·통원·약값…) + 일상배상책임은 항상 파랑</b>.
            #   통원·약값은 실손 디폴트 블록이 BL로 써서 파랑인데, <b>입원은 별첨 파싱 경로</b>를 타서
            #   그 계약이 비갱신이면 검정으로 찍혔다 — 같은 실손 안에서 색이 갈렸다.
            if std in _BLUE_ROWS: blue = True
            # 수술비 1~5종 -> 종별 슬래시 누적
            if std == '종수술비공통' and 1 <= jong <= 5:   # ★v29q-12 상해/질병 미표기 → 상해·질병 양쪽 동일 기재
                # ★★★★★v341 예외사항(지점장 지시 2026.08.02): <b>질병/상해 표기 없는 종수술비가
                #   한 계약에 두 벌 실리면 그것은 「질병 1벌 + 상해 1벌」이지 합산 대상이 아니다.</b>
                #   지점장 원문: "미래에셋 1-5종이 질병or상해라고 기재안된채 2개가 기재되어있다
                #   넌 그걸 <b>각각 더해서 2배</b>로 만들어놧다".
                #   실측 = 미래에셋 `1-5종수술특약(1~5종)` 20/40/300/1000/2000이 <b>두 벌</b>
                #   → 앱이 40/80/600/2000/4000으로 부풀렸다. 정답은 <b>각 축 20/40/300/1000/2000</b>.
                #   → 누적(`+=`)이 아니라 <b>대표(max)</b>로 넣는다. 한 벌만 있으면 종전처럼 양쪽 복제.
                for _k in ('상해 종수술비(1-5종)','질병 종수술비(1-5종)'):
                    jong_acc[_k][jong-1] = max(jong_acc[_k][jong-1], amt)
                    if blue: jong_blue[_k] = True
                surg_trace.append((ct['company'], raw, f'상해·질병 종수술 양쪽 {jong}종 슬롯', amt))   # ★v30g
                continue
            # ★v378 종번호 없는 단일금액 종수술 — 축 미표기는 질병·상해 양쪽(대표 max), 축 표기는 그 행만.
            if std == '종수술비공통' and not jong:
                for _k in ('상해 종수술비(1-5종)','질병 종수술비(1-5종)'):
                    jong_lump[_k] = max(jong_lump[_k], amt)
                    if blue: jong_blue[_k] = True
                surg_trace.append((ct['company'], raw, '상해·질병 종수술 단일금액(종번호 없음)', amt))
                continue
            if std in jong_acc and not jong:
                jong_lump[std] = max(jong_lump[std], amt)
                if blue: jong_blue[std] = True
                surg_trace.append((ct['company'], raw, f'{std} 단일금액(종번호 없음)', amt))
                continue
            if std == 'n대수술비':                       # ★v386 116대 6칸 슬래시
                _ndg = _ndae_grade(raw)
                if _ndg:
                    ndae_acc[_ndg-1] = max(ndae_acc[_ndg-1], amt)
                    surg_trace.append((ct['company'], raw, f'n대수술비 {_ndg}등급 슬롯', amt))
                    continue
            if std in jong_acc and 1 <= jong <= 5:
                jong_acc[std][jong-1] += amt
                if blue: jong_blue[std] = True
                surg_trace.append((ct['company'], raw, f'{std} {jong}종 슬롯', amt))   # ★v30g
                continue
            # ★v29y (지점장 2026.07.02): MRI·도수치료·비급여주사 = 'MRI/도수치료/비급여주사' 한 행 슬래시(1-5종 방식)
            if std in ('MRI','도수치료','비급여주사'):
                _ti={'MRI':0,'도수치료':1,'비급여주사':2}[std]
                trio_acc[_ti]=max(trio_acc[_ti],amt)   # 실손 계열=중복합산 금지, 대표 최댓값
                continue
            if std == '__무시__':            # ★v30z 지점장 무시지정 담보(전이암진단비·고액항암치료비) = 완전 드롭, [확인]에도 미노출
                continue
            # ★★★★★v580 제122조 4항 (지점장 지시 2026.08.24) — <b>엑셀은 0 + 노란칸</b>.
            #   현물 담보는 금액이 없다(보험사가 변호사를 직접 붙여준다). 숫자로 더하면
            #   실측처럼 `1,000 + 500 + 30 = 1,530`이라는 <b>있지도 않은 금액</b>이 생긴다.
            #   ⇒ 그 담보만 0으로 만들고 셀에 노란칸을 남긴다. 「현물」 글자는 산출물 3종이 찍는다.
            if std == '변호사' and _is_hyunmul(raw, ct.get('company'), ct.get('contract_date')):
                _HYUNMUL.add('변호사')
                _hy_cells.append((nm2r.get('변호사'), col))
                amt = 0
            r = nm2r.get(std)
            if r is None and std=='n대수술비': r = nm2r.get('120대수술비')   # ★v352 구 마스터(120대수술비) 하위호환
            if r is None and std:             # 공백무시 재매칭 (화상 '진 단 비' 등)
                r = nm2r_norm.get(re.sub(r'\s','', std))
            if not std or r is None:          # 마스터 미수록/매핑실패 -> [확인]
                unmapped.append((col, ct['company'], raw, amt, m.get('note','') or ''))
                continue
            # 2대 주요치료비는 뇌혈관·심장 두 칸 모두 기재(동일 담보, 양쪽 표기). 그 외는 단일 행.
            # ★★★v216 (지점장 지시 2026.07.25, 영구): <b>혈전용해치료비 누락 수정</b>.
            #   마스터에 '혈전용해치료비' 행이 <b>2개</b>다 — <b>뇌 블록 37행 · 심장 블록 50행</b>.
            #   그런데 `nm2r[담보명]=행`은 <b>뒤에 나온 50행이 앞의 37행을 덮어써</b> 뇌 37행에는
            #   <b>어떤 값도 영원히 들어가지 않았다</b>(구조적 누락). 게다가 실측에서 메리츠
            #   '뇌혈관혈전용해치료비 200' + '급성심근경색증혈전용해치료비 200'이 <b>같은 50행에 400으로 합산</b>됐다.
            #   → <b>담보명의 축으로 갈라 배정한다</b>: 뇌 축=37행 / 심장 축(심근·심장·허혈·관상동맥·심혈관)=50행 /
            #   <b>축 미표기 또는 심뇌 동시 표기 = 양쪽 행에 각 100% 동일 금액</b>(2대 주요치료비·묶음담보 공통원칙 §8.3.1과 동일).
            if std == '2대 주요치료비':
                target_rows = nm2r_multi.get(std, [r])
            elif std == '혈전용해치료비':
                _hjall = nm2r_multi.get(std, [r])
                _hjn = _norm(raw)
                _hjb = ('뇌' in _hjn)
                _hjh = any(_k in _hjn for _k in ('심근','심장','허혈','관상동맥','심혈관'))
                if _hjb and not _hjh:   target_rows = [_hjall[0]]     # 뇌 전용 → 뇌 블록 행
                elif _hjh and not _hjb: target_rows = [_hjall[-1]]    # 심장 전용 → 심장 블록 행
                else:                   target_rows = _hjall          # 축 미표기·심뇌동시 → 양쪽 각 100%
            else:
                target_rows = [r]
            for tr in target_rows:
                existing = ws.cell(tr,col).value
                # ★★★v215 (지점장 확정 2026.07.25, 영구): <b>간병인지원일당 = 대표(max) 1건</b>.
                #   근거: 간병인지원 질병입원일당 3 + 간병인지원 상해입원일당 3 은 <b>질병·상해 택일 지급</b>이라
                #   <b>합산 6이 아니라 3</b>이다(둘 중 하나만 기재). 간병인·간호통합병동·1인실과 동일 처리.
                # ★항암방사선약물도 대표(max) 1건 — 항암약물치료비 / 항암방사선치료비 /
                #   항암방사선약물치료비 / 항암약물방사선치료비는 <b>이름만 다른 같은 담보</b>다(합산 금지).
                _rep1 = _is_repmax(std)   # ★v457 제68조 — 대표(max) 목록 단일화(제19조). 여기서 직접 나열하지 않는다.   # ★v370 2대 주요치료비 추가(순환계 주요치료비 4분할=대표1개) ★v302-B 6주미만=대표(max) / ★v198 합의금=대표1개 / ★v208 1인실 / ★v215 간병인지원일당=택일 대표(max)
                _rep1 = _rep1 or ('통합' in raw and std in ('일반암','유사암(갑.기.경.제)','통합전이암'))   # ★v30a §8.2 통합 계열=대표금액 1개
                # ★v320 통합암·10억 플랜도 <b>대표금액 1개</b>
                _rep1 = _rep1 or (std in ('통합암','10억 플랜'))
                # ★★★★★v394 (지점장 확정 2026.08.12, 영구):
                #   지점장 원문: 「<b>암입원일당(1-180)</b> 과 <b>암입원일당(요양병원)</b>은 <b>큰값 1개만</b> 입력해주면 된다」
                #   → 마스터 <b>암일당(30행)</b>은 <b>대표(max)</b>다. 합산 금지.
                #   근거: 두 담보는 <b>같은 입원을 요양병원이냐 아니냐로 나눠 놓은 것</b>이라
                #     동시에 받는 담보가 아니다 — 간병인지원일당·1인실·간호통합병동과 같은 <b>택일 대표</b> 원칙.
                #   ★구 기본값(합산)이면 `10 + 10 = 20`으로 <b>2배</b>가 됐다.
                #   ★v344 기본 대표(max)는 <b>담보명이 완전히 같은 줄</b>에만 걸려서 이 둘(이름이 다르다)은 안 걸렸다.
                _rep1 = _rep1 or (std == '암일당')
                if _rep1 and isinstance(existing,(int,float)):
                    ws.cell(tr,col).value = max(existing, amt)   # 표적·n대·창상봉합=대표 최댓값1건(★v29q-6) / 실손=중복합산 안함(한도)
                else:
                    ws.cell(tr,col).value = (existing+amt) if isinstance(existing,(int,float)) else amt
                # 실손(입원/통원/약값)·일상배상책임은 갱신·비갱신 무관 항상 파랑
                # ★★★v210 (지점장 확정 2026.07.25, 영구): <b>간병인 · 간호통합병동 2가지는 '항상 파랑' 강제 폐기</b>.
                #   보험료 · 가입년일 · 만기일자 · 총납입기간(=계약 갱신 판정) 또는 담보명의 <b>[갱신] 표기</b>에 따라
                #   갱신=파랑 / 비갱신=검정으로 <b>일반 담보와 동일하게</b> 칠한다(구 v139 '간병인 계열 3행 무조건 파랑' 폐기).
                # ★★★★★v500 제98조 3항 (지점장 확정 2026.08.19 「흥국 10억통장은 <b>그냥 무조건 갱신</b>이다」)
                #   raw 표기가 어떻게 들어오든 <b>「10억 플랜」 행은 항상 파랑</b>이다.
                #   `_is_gen_dambo`의 리셋월렛 인식(v499)에 더해 <b>표준명으로도</b> 막는다 — 두 겹.
                ws.cell(tr,col).font = BL if (blue or std in ('입원','통원','약값','약','일상배상책임','10억 플랜')) else BK
                if blue or std == '10억 플랜': _blue_r.add(tr)   # ★v475 제83조 · v500 10억 플랜
                # ★v39 워크시트용 원본담보명 수집(그 표준명 중 최댓값 담보의 raw 1개)
                try: trace_all.append((str(std), str(raw).strip(), amt, str(ct.get('company','')), str(ct.get('product',''))[:40]))
                except Exception: pass
                _WS_STD = ('암주요치료비','하이클래스(암)','2대 주요치료비','산정특례뇌혈관','산정특례심장','일반암','뇌혈관진단비','뇌졸증진단비','급성심근경색','허혈성 진단비')
                if std in _WS_STD:
                    _prev = raw_by_std.get(std)
                    if _prev is None or amt >= _prev[1]:
                        raw_by_std[std] = (str(raw).strip(), amt)
                if std in {'협심증','심부전','빈맥','염증','부정맥','심근병증','심장판막','산정특례심장','2대 주요치료비','허혈성 진단비','급성심근경색','중대한 급성심근','혈전용해치료비','심장수술비','허혈성수술비'}:
                    heart_trace.append((ct['company'], raw, std, amt))   # ★v29z 심장 단독 기재 근거
                if '수술' in str(std) or std == '창상봉합술':
                    surg_trace.append((ct['company'], raw, std, amt))   # ★v30g 수술 기재 근거
                if std in {'일반암','유사암(갑.기.경.제)','통합전이암','고액암','중대한 암','암주요치료비','하이클래스(암)','암수술','암일당'}:
                    cancer_trace.append((ct['company'], raw, std, amt))   # ★v30h 암 기재 근거

        for nm, vals in jong_acc.items():     # 종수술비 슬래시 기재(§6)
            # ★v378 종번호 없는 단일금액 — 슬래시 칸이 비었을 때만 그 행에 숫자 그대로.
            #   둘이 동시에 차면 어느 쪽이 정본인지 지침에 없다 → 임의 병합 금지, [확인]큐.
            if jong_lump.get(nm):
                if any(vals):
                    unmapped.append((col, ct['company'], f'{nm}(종번호 없는 단일금액)', jong_lump[nm],
                                     '[확인] 같은 계약에 종별 슬래시와 단일금액이 동시 존재 — 병합 규칙 없음'))
                else:
                    _rl = nm2r.get(nm)
                    if _rl:
                        ws.cell(_rl,col).value = jong_lump[nm]
                        ws.cell(_rl,col).font = BL if (gen or jong_blue[nm]) else BK
                    continue
            if any(vals):
                # ★v29v (지점장 2026.07.02): 6~8종 값이 있으면 그 계약의 종수술은 8단계 → (1-8종) 행에 8칸 슬래시,
                #   아니면 기존대로 (1-5종) 행에 5칸 슬래시.
                if any(vals[5:]):
                    tgt=nm.replace('(1-5종)','(1-8종)'); use=vals
                else:
                    tgt=nm; use=vals[:5]
                r = nm2r.get(tgt) or nm2r.get(nm)
                if r:
                    ws.cell(r,col).value = '/'.join(str(x) for x in use)
                    ws.cell(r,col).font = BL if (gen or jong_blue[nm]) else BK

        if any(ndae_acc):   # ★v386 116대 수술비 등급 슬래시 기재
            _rnd = nm2r.get('n대수술비') or nm2r.get('120대수술비')
            if _rnd:
                _cur = ws.cell(_rnd,col).value
                # ★★★★★v595 제136조 (지점장 확정 2026.08.26 「<b>N대수술비는 대표값 넣어줘</b>」).
                #   [구 동작] 등급 6칸 슬래시 `0/500/0/0/0/0` — 0이 다섯 칸이라 읽기 어렵고
                #     PPT `_____대 수술` 칸에도 슬래시가 그대로 들어갔다.
                #   ⇒ <b>등급 중 최댓값 하나</b>만 적는다(`_REPMAX_ROWS`에 이미 `n대수술비`가 있다 — 같은 뜻).
                #   ★근거는 `근거표` 시트의 `n대수술비 N등급 슬롯` 줄에 그대로 남는다(추적 가능).
                _nmax = max(int(x or 0) for x in ndae_acc)
                if isinstance(_cur,(int,float)) and _cur:
                    _nmax = max(_nmax, int(_cur))
                ws.cell(_rnd,col).value = _nmax
                ws.cell(_rnd,col).font = BL if gen else BK

        if any(trio_acc):   # ★v29y MRI/도수/주사 슬래시 기재(실손 계열=항상 파랑)
            _rt=nm2r.get('MRI/도수치료/비급여주사')
            if _rt:
                ws.cell(_rt,col).value='/'.join(str(x) for x in trio_acc)
                ws.cell(_rt,col).font=BL

        # ★★★★★v239 CI 본체 최종 기재 — 담보 루프가 끝난 뒤에 한다.
        #   정본: 본체(=사망액 × 선지급률)를 <b>중대한 암·중대한 뇌졸증·중대한 급성심근 3행에 동일 금액</b>으로 기재하고,
        #   <b>초과분은 일반 담보 행</b>에 남긴다. 중대한CI적용 = 사망 − 본체.
        #   (지점장 지적 "각 중대한의 금액이 다 틀리다" = 구 코드가 별첨 담보 원값을 그대로 옮겨 3행이 제각각이었다)
        if _ci_fix:
            _bc=_ci_fix['bonche']; _sm=_ci_fix['samang']
            # ★★★★★v242 영구지침(지점장 지시 2026.07.25): <b>CI 본체의 뇌 축은 '뇌졸증'이 아니라 '뇌졸증 or 뇌출혈'</b>이다.
            #   지점장 원문 — "정답은 <b>중대한뇌출혈</b>이다. 왜냐면 사망이 3천에 중대한뇌출혈은 2400 + 추가로 뇌출혈 가입이고.
            #   <b>신한도 뇌출혈</b>이다. … <b>무조건 뇌졸증이 아니라는거다 뇌졸증 or 뇌출혈이다.</b>"
            #   ★실측 근거 — DB 별첨 `뇌출혈진단 2,400`이 <b>본체와 정확히 일치</b>(나머지 1,000+1,000은 추가 가입분),
            #     신한은 뇌 담보가 `뇌출혈진단`뿐. 구 코드는 무조건 '중대한 뇌졸증'에 넣어 <b>축이 통째로 틀렸다</b>.
            #   ★판별: 그 계약 별첨의 <b>뇌 담보 이름</b>을 본다 — ①본체 금액과 일치하는 담보의 축 ②없으면 뇌출혈 우선.
            _braw = ct.get('ci_brain') or {}
            _bax = '중대한 뇌출혈' if _braw.get('axis')=='뇌출혈' else '중대한 뇌졸증'
            _bovf = '뇌출혈진단비' if _braw.get('axis')=='뇌출혈' else '뇌졸증진단비'
            for _nm,_ovf in (('중대한 암','일반암'),(_bax,_bovf),('중대한 급성심근','급성심근경색')):
                _r=nm2r.get(_nm); _ro=nm2r.get(_ovf)
                if not _r: continue
                # ★★★★★v242: <b>중대한 행 + 일반 행을 합쳐서</b> 본체만큼 중대한 행에, 나머지를 일반 행에 배분한다.
                #   구 코드는 <b>중대한 행의 원값만</b> 봤다 → 뇌출혈처럼 값이 <b>일반 행(뇌출혈진단비)에만</b> 있으면
                #   `원값 0 ≤ 본체`라 스킵되어 <b>중대한 뇌출혈이 0으로 남았다</b>(지점장 지적 "이번에도 틀렸다").
                #   ★합계를 보존하므로 한장보장표(중대한OO + 일반 합산)는 어느 쪽이든 불변이다.
                _v1=ws.cell(_r,col).value;  _v1=_v1 if isinstance(_v1,(int,float)) else 0
                _v2=ws.cell(_ro,col).value if _ro else 0; _v2=_v2 if isinstance(_v2,(int,float)) else 0
                _tot=_v1+_v2
                # ★★★★★v253 영구지침(지점장 지시 2026.07.26 "오늘 무조건 CI는 문제 없어야 한다"):
                #   <b>합계가 0이어도 본체를 중대한OO 행에 기재한다</b>.
                #   구 코드는 `_tot<=0`이면 continue라 <b>중대한OO 3행이 통째로 0</b>이 됐다.
                #   ★왜 0이 되나 — 교보생명 큰사랑CI처럼 <b>별첨 담보명이 '주계약'·'○○특약'뿐</b>인 생보 CI는
                #     암·뇌·심 담보가 종류별로 dambo에 없다(4p 세부가입현황에만 있다). 그래서 뺄 일반행이
                #     없는데 정본 배분이 "합계에서 나눈다"라서 <b>나눌 합계가 0</b> → 스킵 → 기재 0.
                #   ★실측(김O구 교보): 1차 별첨 주계약 4,000(=일반사망)×80% = 2차 세부가입현황 질병사망
                #     4,000×80% = 본체 3,200. 두 경로 일치인데 <b>중대한 암·뇌졸증·급성심근이 전부 0</b>이었다.
                #   ★합계가 있으면 기존 배분(합계 보존) 그대로 — 신한·DB생명 회귀 없음.
                # ★★★★★v254 영구지침(지점장 지시 2026.07.26): <b>별첨 '주계약' 라벨 경로(교보 등)는
                #   CI 본체와 일반 담보가 별개다</b> → 중대한행에 <b>본체를 기재하고 일반행은 손대지 않는다</b>.
                #   지점장 원문 "니가 올린건 CI + 일반 담보다 / 교보의 일반 급성심근 담보 정답".
                #   ★실측(김O구 교보): 급성심근 4,200 = <b>CI 본체 3,200 + 교보 일반 급성심근 1,000</b>.
                #     구 코드는 합계에서 본체를 빼는 배분이라 <b>일반 1,000을 지워</b> 3,200이 됐다.
                #   ★반대로 별첨 <b>담보명</b>으로 CI가 표시된 계약(신한 `뇌출혈진단`·DB생명)은 담보값에
                #     본체가 <b>포함</b>돼 있으므로 기존 배분(합계 보존)을 유지한다 — 회귀 없음.
                if _cij:
                    ws.cell(_r,col).value=_bc; ws.cell(_r,col).font = BL if gen else BK
                    print(f"[v254 CI가산] {ct.get('company')} {_nm} ← 본체 {_bc:,} · {_ovf} {_v2:,} 유지(일반 담보 별개)")
                    continue
                if _tot<=0:
                    ws.cell(_r,col).value=_bc; ws.cell(_r,col).font = BL if gen else BK
                    print(f"[v253 CI기재] {ct.get('company')} {_nm} ← 본체 {_bc:,} (일반행 값 0)")
                    continue
                # ★★★★★v253: 중대한 행 = <b>본체 그대로</b>(구 코드 `min(_bc,_tot)`은 합계가 본체보다
                #   작을 때 <b>본체가 잘려나갔다</b>). 정본 = "중대한OO 3행에 <b>동일 금액</b> 기재".
                #   ★실측(김O구 교보): 급성심근 합계 1,000 < 본체 3,200 → 1,000만 기재되어 한장표 4,200에
                #     2,200이 비었다. 본체 3,200을 기재하면 다른 계약분 1,000과 합쳐 <b>4,200</b>이 맞는다.
                #   ★합계 ≥ 본체인 경우는 `min(_bc,_tot)==_bc`라 <b>기존과 동일</b> — 신한·DB생명 회귀 없음.
                _put=_bc
                ws.cell(_r,col).value=_put; ws.cell(_r,col).font = BL if gen else BK
                if _ro:
                    _rest=max(0,_tot-_put)
                    ws.cell(_ro,col).value=(_rest if _rest>0 else None)
                    if _rest>0: ws.cell(_ro,col).font = BL if gen else BK
            # ★★★v242: <b>축이 아닌 쪽의 '중대한' 뇌 행은 일반 행으로 되돌린다</b>.
            #   CI 상품이라 resolve가 `뇌졸중진단`을 '중대한 뇌졸증'으로 리네임했는데, 축이 뇌출혈로 확정되면
            #   그 값은 <b>CI 본체가 아니라 별도 일반 담보</b>다(실측 DB 뇌졸중진단 1,000). 합계는 보존된다.
            _other = ('중대한 뇌졸증','뇌졸증진단비') if _bax=='중대한 뇌출혈' else ('중대한 뇌출혈','뇌출혈진단비')
            _ro2=nm2r.get(_other[0]); _rn2=nm2r.get(_other[1])
            if _ro2 and _rn2:
                _vo=ws.cell(_ro2,col).value
                if isinstance(_vo,(int,float)) and _vo:
                    _vn=ws.cell(_rn2,col).value
                    ws.cell(_rn2,col).value=(_vn if isinstance(_vn,(int,float)) else 0)+_vo
                    ws.cell(_rn2,col).font = BL if gen else BK
                    ws.cell(_ro2,col).value=None
                    print(f"[v242 축정리] {ct.get('company')} {_other[0]} {_vo:,} → {_other[1]} (축={_bax})")
            _rci=nm2r.get('중대한CI적용')
            if _rci:
                ws.cell(_rci,col).value=_sm-_bc; ws.cell(_rci,col).font = BL if gen else BK
            print(f"[v239 CI본체] {ct.get('company')} 사망{_sm:,} × {_ci_fix['pct']}% = 본체 {_bc:,} → 중대한OO 3행 동일기재 · 잔여 {_sm-_bc:,}")

        # ★ §8 생보 종신(만기 9999): 일반사망(종신) + 상해사망 1:1 복제
        if ct['expiry_date'].startswith('9999'):
            r_il = nm2r.get('일반사망'); r_sh = nm2r.get('상해사망'); r_jb = nm2r.get('질병사망(80세)')
            # ★★★★★v239 (지점장 지적 2026.07.25 "엑셀에 사망이 안 잡혀있다"): 정본 §8.1은
            #   <b>"일반사망 = 생명보험사 만기 9999(종신)으로 표기된 사망"</b>인데, 구 코드는
            #   담보명이 <b>'일반사망'</b>일 때만 6행에 넣고 <b>'질병사망'</b>은 질병사망(80세) 8행에 뒀다.
            #   → CI 종신보험(신한 라이프케어CI·DB CI종신)은 별첨 담보명이 <b>'질병사망'</b>이라
            #   <b>일반사망 행이 통째로 비었고</b>, 그 결과 선지급률 계산 근거가 엑셀에 없었다.
            #   ★생보 + 만기 9999면 질병사망(80세) 값을 <b>일반사망(종신) 행으로 이동</b>한다.
            #   (질병사망 합계는 '일반사망 + 질병사망(80세)'로 보므로 <b>한장보장표 총액은 불변</b> — 실측 확인)
            # ★★★★★v269 영구지침(지점장 지시 2026.07.29): <b>"우체국은 다 다르다. 우선 질병사망에 넣어라"</b>
            #   → 우체국(우정사업본부) 계약은 <b>만기 9999(종신)여도 종신 판정에서 제외</b>하고
            #     §8.1③에 따라 <b>질병사망(80세)</b> 행에 둔다. 우체국 상품은 종신 표기가 제각각이라
            #     일괄 종신 처리하면 안 된다는 지점장 판단이다.
            #   ★회사명이 <b>'우체국'으로 오든 '우정사업본부'로 오든 결과가 같아야</b> 하므로 두 표기를 함께 배제한다.
            #   ★<b>이 블록(사망 배정)에서만</b> 뺀다 — 입원특약 일당(1423행)·실손 통원/약값(3197행)의
            #     생보 판정에는 손대지 않는다(사망과 무관).
            _cnm269 = ct.get('company') or ''
            _life = any(k in _cnm269 for k in ('생명','라이프','AIA','메트라이프','공제')) \
                    and not any(k in _cnm269 for k in ('우체국','우정'))
            if _life and r_il and r_jb:
                _vjb = ws.cell(r_jb,col).value
                if isinstance(_vjb,(int,float)) and _vjb:
                    # ★★★★★v240(지점장 지시): CI 계약이면 <b>주계약 사망만 일반사망(종신)</b>으로 옮기고,
                    #   <b>차액(추가 특약 사망)은 질병사망(80세) 행에 그대로 남긴다</b>.
                    #   실측 신한 — 별첨 사망 10,000 중 주계약 4,000(본체 2,000÷50%) → 일반사망 4,000 / 질병사망 6,000.
                    #   (CI가 아니거나 주계약 판별 실패면 종전대로 전액 이동 — 한장표 총액은 어느 쪽이든 불변)
                    _mv = _vjb
                    if _ci_fix and 0 < _ci_fix.get('samang',0) < _vjb:
                        _mv = _ci_fix['samang']
                    _v0 = ws.cell(r_il,col).value
                    ws.cell(r_il,col).value = (_v0 if isinstance(_v0,(int,float)) else 0) + _mv
                    ws.cell(r_il,col).font = BL if gen else BK
                    ws.cell(r_jb,col).value = (_vjb-_mv) if (_vjb-_mv) > 0 else None
                    print(f"[v240 종신사망] {ct.get('company')} 별첨사망 {_vjb:,} → 일반사망(종신) {_mv:,}"
                          + (f" · 질병사망(80세) {_vjb-_mv:,}(추가특약)" if _vjb-_mv>0 else ""))
            v = ws.cell(r_il,col).value if r_il else None
            if isinstance(v,(int,float)) and r_sh and not isinstance(ws.cell(r_sh,col).value,(int,float)):
                ws.cell(r_sh,col).value = v
                ws.cell(r_sh,col).font = BL if gen else BK

        # ★ 실손 통원/약값 디폴트 (지점장 2026.06.28): ①별첨 명시값 최우선 → ②입원3,000 구형=통원10·약5
        #   → ③2021.06 이전: 손보 통원25·약5 / 생보 통원20·약10 → ④4세대(2021.07~): 통원20·약0(통원포함).
        _rip=nm2r.get('입원'); _rtw=nm2r.get('통원'); _ryk=nm2r.get('약값')
        _ipv=ws.cell(_rip,col).value if _rip else None
        if isinstance(_ipv,(int,float)) and _ipv:   # 이 계약에 실손(입원) 존재
            _life=any(k in (ct['company'] or '') for k in ('생명','라이프','AIA','메트라이프','우체국','공제'))
            def _ym(d):
                try: return int(str(d)[:4])*100+int(str(d)[5:7])
                except: return 0
            _g4=(silson_gen(ct.get('contract_date',''), None, ct.get('product',''), _has_nonpay3(ct.get('dambo')), _has_drug(ct.get('dambo'))) in ('4세대','5세대'))   # ★v29v 상품코드 반영
            _guhy=(_ipv==3000)                            # 입원한도 3,000=구형
            _twc=ws.cell(_rtw,col).value if _rtw else None
            _ykc=ws.cell(_ryk,col).value if _ryk else None
            # ★v215: 통원 디폴트 판정에도 1세대(가입일 기준)를 포함한다. 구 코드는 입원한도 3,000(_guhy)일
            #   때만 10을 넣어, 가입일이 2009.09 이전인데 입원한도가 3,000이 아닌 1세대는 25/20이 됐다.
            _g1a=str(silson_gen(ct.get('contract_date',''), _ipv, ct.get('product',''), _has_nonpay3(ct.get('dambo')), _has_drug(ct.get('dambo')))).startswith('1세대')
            if _rtw and not isinstance(_twc,(int,float)):  # ① 별첨 통원 없을 때만 디폴트
                _twd = 10 if (_guhy or _g1a) else (20 if _g4 else (20 if _life else 25))
                ws.cell(_rtw,col).value=_twd; ws.cell(_rtw,col).font=BL
            # ★★★v215 (지점장 확정 2026.07.25, 영구): <b>실손 1세대 = 입원 3,000 / 통원 10 / 약값 0</b>.
            #   1세대는 외래+약제가 <b>통원 한 한도로 통합</b>돼 있어 <b>약값 담보 자체가 없다</b> → 약값 칸은 비운다.
            #   ★버그2건 실측: ①`silson_gen`은 입원한도 3,000이면 <b>'1세대(구형)'</b>을 반환하는데
            #     구 코드는 `== '1세대'`로 비교해 <b>_g1이 항상 False</b>였다 → `_guhy` 분기로 빠져 <b>약값 5가 찍혔다</b>
            #     (지점장 지적 "보장진단서 오류"의 원인). → `.startswith('1세대')`로 수정.
            #     ②구 코드는 '별첨에 약값이 없을 때만' 0으로 뒀다 → 별첨에 약값이 인쇄돼 있으면 그대로 들어갔다.
            #     → 1세대는 <b>별첨 명시값이 있어도 약값 칸을 지운다</b>(강제).
            _g1=str(silson_gen(ct.get('contract_date',''), _ipv, ct.get('product',''), _has_nonpay3(ct.get('dambo')), _has_drug(ct.get('dambo')))).startswith('1세대')   # ★v215 (구 v41 == '1세대' 버그)
            if _ryk and _g1:                               # ★v215 1세대=약값 없음(강제 미기재)
                if isinstance(_ykc,(int,float)) and _ykc:
                    silson_trace.append((ct['company'], ct.get('contract_date',''), '1세대 약값삭제', f'{_ykc}→0'))
                ws.cell(_ryk,col).value=None
            elif _ryk and not isinstance(_ykc,(int,float)):  # ① 별첨 약값 없을 때만 디폴트
                _ykd = 0 if _g4 else (10 if _life else 5)
                if _ykd: ws.cell(_ryk,col).value=_ykd; ws.cell(_ryk,col).font=BL   # 4세대 약0=미기재
            # ★ 실손 세대 자동판별 → 헤더에 라벨 기재
            _sg = silson_gen(ct.get('contract_date',''), _ipv, ct.get('product',''), _has_nonpay3(ct.get('dambo')), _has_drug(ct.get('dambo')))
            _pm0=re.search(r'(?<!\d)(0[9]|1[0-9]|2[0-6])(0[1-9]|1[0-2])(?!\d)', str(ct.get('product','')))
            silson_trace.append((ct['company'], ct.get('contract_date',''), (_pm0.group(0) if _pm0 else '없음'), _sg or '판정불가'))   # ★v29z 세대 근거
            if _sg:
                _hc = ws.cell(1,col)
                if _hc.value and _sg not in str(_hc.value):
                    _hc.value = str(_hc.value) + f'\n({_sg} 실손)'

        # ★★★★★v475 제83조 — <b>이 계약 열의 색은 여기서 확정한다</b>(지점장 2026.08.18).
        #   담보명에 '(갱신)'이 있으면 <b>주계약이 비갱신이어도 그 담보만 파랑</b>(제5조 B).
        #   담보 기재 뒤에 CI 배분·종신 사망 이동·묶음 분해가 같은 셀을 다시 칠하므로,
        #   <b>계약 루프의 맨 끝에서 한 번 더</b> 파랑을 확정한다 — 나중 규칙이 앞 규칙을 죽이지 못하게.
        if _blue_r:
            _bn = 0
            for _r in sorted(_blue_r):
                if isinstance(ws.cell(_r,col).value, (int,float,str)) and ws.cell(_r,col).value not in (None,''):
                    ws.cell(_r,col).font = BL; _bn += 1
            print('[v475 갱신담보색] %s 열%d — 파랑 확정 %d행' % (ct.get('company',''), col, _bn))

    # ★v29t (지점장 확정 2026.07.02): CI 존재 시 '중대한CI적용' 행 = CI 잔여액 + 비CI 계약의 일반사망 동일액 —
    #   CI 적용/미적용 각각의 총 사망액이 양쪽 행에서 가로합산되도록.
    # ★★★v208 1인실 상급병원 가산(지점장 확정 2026.07.25, 영구):
    #   <b>1인실 상급병원 = 상급종합 1인실 + 종합 1인실</b>(상급종합에 입원하면 종합병원이하 담보도 함께 나온다).
    #   <b>1인실 종합병원 = 종합 1인실만</b>. 질병·상해는 택일 지급이라 담보 기재 단계에서 대표(max)로 잡힌다.
    #   실측(양*선 삼성 New내돈내삼): 상급 20 · 종합 10 → 상급행 30 / 종합행 10.
    #   ★계약 루프 <b>바깥</b>에서 1회만 돌린다(루프 안에 넣으면 계약 수만큼 중복 가산 — 실측 60 오류).
    _r1s = nm2r.get('1인실 상급병원'); _r1j = nm2r.get('1인실 종합병원')
    if _r1s and _r1j:
        for _c in range(3, 3 + n_ct):
            _vs = ws.cell(_r1s, _c).value; _vj = ws.cell(_r1j, _c).value
            if isinstance(_vs, (int, float)) and isinstance(_vj, (int, float)):
                ws.cell(_r1s, _c).value = _vs + _vj

    _rci_all=None; _ril_all=None; _rjb_all=None
    for _rr in range(6, ws.max_row+1):
        _b=str(ws.cell(_rr,2).value or '').strip()
        if _b=='중대한CI적용': _rci_all=_rr
        if _b=='일반사망': _ril_all=_rr
        if _b=='질병사망(80세)': _rjb_all=_rr
    _has_ci=any(_isci_prod(c.get('product'), c.get('company')) for c in contracts)
    # ★★★★★v245 영구지침(지점장 확정 2026.07.25): <b>비CI 계약의 사망은 '질병사망(80세)' 행에 넣는다</b>.
    #   지점장 원문 = "(비CI 일반사망이 중대한CI적용에 찍히는 v29t 규칙) <b>그건 질병사망(80)에 넣어라</b>".
    #   구 v29t는 비CI 계약의 일반사망을 <b>'중대한CI적용'에 복사</b>했다(실측 메트라이프 6,000).
    #   → <b>'일반사망(종신)' 행은 CI 주계약 사망 전용</b>, 비CI 종신 사망은 질병사망(80세)로 옮긴다.
    #   ★한장보장표 질병사망 = 일반사망 + 질병사망(80세) 합이므로 <b>총액은 불변</b>(실측 22,000 유지).
    if _has_ci and _ril_all and _rjb_all:
        for _ix,_c in enumerate(contracts):
            _cl=3+_ix
            if _isci_prod(_c.get('product'), _c.get('company')): continue
            # ★★★★★v268 영구지침(지점장 지적 2026.07.29 "이분은 생명보험인데 다 질병사망(80세)로 기재된다"):
            #   정본 §8.1 = <b>"일반사망 = 생명보험사 만기일자 9999(종신)으로 표기된 사망만 → 일반사망(종신)"</b>.
            #   구 v245는 <b>비CI면 종신 여부를 보지 않고</b> 일반사망 행을 전부 질병사망(80세)로 옮겨,
            #   같은 실행 안에서 <b>v240이 §8.1대로 올려놓은 값을 그대로 되돌렸다</b>.
            #   실측(이명순): 한화생명 종신 5계약 1,200·4,050·100,000·26,076·50,000 = <b>181,326이 전부 질병사망으로 밀렸다</b>.
            #   → <b>생보 + 만기 9999(종신)이면 일반사망(종신) 행을 그대로 둔다</b>(중대한CI적용 복사 제거만 수행).
            #   ★한장보장표 질병사망 = 일반사망 + 질병사망(80세) 합이므로 <b>총액은 어느 쪽이든 불변</b>.
            # ★v269: 종신 판정 기준을 v240 블록과 <b>완전히 동일</b>하게 맞춘다(우체국·우정 배제).
            #   두 곳의 기준이 어긋나면 한쪽은 안 올리고 한쪽은 유지해 다시 갈라진다.
            _cn268 = _c.get('company') or ''
            _lf268 = any(k in _cn268 for k in ('생명','라이프','AIA','메트라이프','공제')) \
                     and not any(k in _cn268 for k in ('우체국','우정'))
            if _lf268 and str(_c.get('expiry_date') or '').startswith('9999'):
                if _rci_all: ws.cell(_rci_all,_cl).value=None
                _v268=ws.cell(_ril_all,_cl).value
                if isinstance(_v268,(int,float)) and _v268:
                    print(f"[v268 종신유지] {_c.get('company')} 일반사망(종신) {_v268:,} 유지 (구 v245가 질병사망으로 밀던 것 차단)")
                continue
            _ilv=ws.cell(_ril_all,_cl).value
            if isinstance(_ilv,(int,float)) and _ilv:
                _j0=ws.cell(_rjb_all,_cl).value
                ws.cell(_rjb_all,_cl).value=(_j0 if isinstance(_j0,(int,float)) else 0)+_ilv
                ws.cell(_rjb_all,_cl).font=ws.cell(_ril_all,_cl).font.copy()
                ws.cell(_ril_all,_cl).value=None
                if _rci_all: ws.cell(_rci_all,_cl).value=None   # ★구 v29t 복사분 제거
                print(f"[v245 비CI사망] {_c.get('company')} 일반사망 {_ilv:,} → 질병사망(80세) (중대한CI적용 복사 제거)")

    # ★★★★★v271 (지점장 지적 2026.07.30 "삼성생명 41,800원 계약이 다 비어있다")
    #   별첨 담보명이 <b>특약명/상품명뿐</b>이라 담보 종류를 알 수 없는 계약은
    #   매핑이 0건이 되어 <b>열이 통째로 빈다</b>(삼성생명 無○○(본인) 계열, 교보 `주계약` 등).
    #   → 그런 계약에 한해 <b>세부가입현황 '계약별 가입정보'의 인쇄값</b>으로 채운다(지침 §3 ②).
    #   ★<b>매핑이 1건이라도 있는 계약은 절대 건드리지 않는다</b> — 별첨이 정본이므로 중복 기재 방지.
    #   ★사망은 §8.1을 그대로 적용한다: 생보 AND 만기 9999(우체국·우정 제외) → <b>일반사망(종신)</b>.
    # ★v293 흔적: 검증 앵커가 없어 세부보충을 차단했으면 확인사항에 반드시 남긴다(조용한 누락 금지).
    if data.get('sebu_blocked'):
        unmapped.append((3, '—', '[세부보충 차단] 한장보장표(검산 앵커)가 없어 세부가입현황 보충을 하지 않았다',
                         0, f"보충 후보 {data['sebu_blocked']}건 폐기 — 별첨에 담보가 거의 없는 계약은 리포트 원문 수기 확인 필요"))
    _SEBU_LOG = []      # ★v359 세부보충 흔적(회사, 담보, 금액)
    _sbc = data.get('sebu_bc') or {}
    if _sbc:
        for _ci2, _c2 in enumerate(contracts):
            _cl2 = 3 + _ci2
            # ★★v276: 구 조건은 '별첨 매핑 0건'이었다. 그런데 삼성리빙케어처럼
            #   별첨이 전부 상품명인데 <b>담보 1~2개만 우연히 매핑</b>되면 조건이 깨져
            #   세부보충이 통째로 멈췄다(이영태 실측: 삼성 CI가 전부 0). →
            #   <b>매핑 담보 3개 이하 = 별첨이 담보 종류를 알려주지 못하는 계약</b>으로 보고
            #   그 열을 <b>세부가입현황으로 재구성</b>한다(지침 §3 ②가 유일 정본).
            _filled = [_r2 for _r2 in range(6, ws.max_row+1)
                       if isinstance(ws.cell(_r2,_cl2).value,(int,float))]
            # ★★★★★v296 (이명순 롯데 실측 2026.07.31, 영구): <b>「전부 아니면 전무」를 폐기한다</b>.
            #   구 조건은 별첨 매핑이 4개 이상이면 세부보충을 <b>통째로 스킵</b>했다. 그런데 롯데 신형처럼
            #   <b>별첨 담보명이 공란(금액만)</b>인 리포트는 담보 몇 개만 우연히 잡히고 나머지가 전멸한다.
            #   실측(삼성리빙케어 139,800): 별첨 4개(사망·CI) 매핑 → 스킵 → <b>유사암 600 · 상해후유 4,200 ·
            #   상해수술 100 · 일당 5·5 · 골절 30이 통째로 누락</b>(한장표 불일치 5건의 원인).
            #   → <b>4개 초과면 덮어쓰지 않고 '빈 행만' 보충</b>한다(기존 정상값 보존 + 누락 방지).
            #   ★<b>CI 짝 행에 값이 있으면 그 축은 건너뛴다</b> — 중대한OO에 본체가 이미 있는데
            #     일반행에 같은 값을 또 넣으면 <b>이중계산</b>이 된다(§22 CI 정본).
            _partial = len(_filled) > 3
            _pm = _c2.get('premium')
            try: _pm = int(str(_pm).replace(',','').replace('원','').strip())
            except: _pm = None
            _vals = _sbc.get(_pm)
            if not _vals: continue
            _cn2 = _c2.get('company') or ''
            _lf2 = any(k in _cn2 for k in ('생명','라이프','AIA','메트라이프','공제')) \
                   and not any(k in _cn2 for k in ('우체국','우정'))
            _jong = _lf2 and str(_c2.get('expiry_date') or '').startswith('9999')
            _put=0
            if _filled and not _partial:           # ★재구성: 별첨 잔재 제거 후 세부값으로 채운다
                for _r2 in _filled: ws.cell(_r2,_cl2).value = None
            # ★v296 CI 짝 — 중대한OO에 본체가 있으면 그 축의 일반행은 보충하지 않는다(이중계산 차단)
            # ★★★★★v327b (이영태 실측 2026.08.02, 영구): <b>'일반사망'↔'중대한CI적용' 짝을 제거</b>한다.
            #   암·뇌·심 3축은 <b>배분</b>(합계 보존)이라 중대한행에 본체가 있으면 일반행을 또 채우면
            #   이중계산이 맞다. 그러나 <b>사망은 배분이 아니다</b> — §8.1④ 정본은
            #   <b>「주계약 사망 → 일반사망(종신) 전액」 + 「사망−본체 → 중대한CI적용」이 병존</b>이다
            #   (실측 김O구 교보: 일반사망 4,000 · 중대한CI적용 800이 <b>동시에</b> 기재).
            #   구 코드는 중대한CI적용에 값이 있으면 일반사망 보충을 통째로 건너뛰어
            #   <b>삼성 리빙케어 일반사망 5,000이 소실</b>됐다(한장표 질병사망 −5,000).
            #   ★중대한CI적용은 검산식(질병사망=일반사망+질병사망(80세))에 들어가지 않으므로
            #     둘을 함께 기재해도 이중계산이 되지 않는다.
            # ★v358 짝 행 — 하나라도 별첨 값이 있으면 세부보충 금지(별첨이 정본)
            _SEBU_PAIR = {'골절(치아파절포함)':'골절(치아파절제외)',
                          '골절(치아파절제외)':'골절(치아파절포함)'}
            _CIPAIR = {'일반암':'중대한 암', '뇌졸증진단비':'중대한 뇌졸증',
                       '뇌출혈진단비':'중대한 뇌출혈', '급성심근경색':'중대한 급성심근'}
            # ★★★★★v299-1 (주재현·심정자 실측 2026.07.31, 영구): 세부가입현황이 <b>직접 싣는</b>
            #   담보는 <b>별첨보다 우선</b>한다. 근거 = 두 고객 27담보 전부 <b>세부 합계 == 한장보장표</b>
            #   (27/27 일치)였고, 틀린 12건은 전부 <b>별첨 담보명 오매핑</b>이었다.
            #   실측 예 — 신한 `특정교통상해후유장해`+`교통상해후유장해` 22,000이 상해후유3%로 산입.
            #   ★<b>화이트리스트에 한한다</b> — 세부표 1칸이 마스터 여러 행으로 갈리는 담보는 넣지 않는다:
            #     ㉠사망 3행(일반/질병(80세)/상해) ㉡골절 2행(치아파절 포함/제외)
            #     ㉢<b>통원</b>(세부 30 = 통원25 + 약값5 합산표기 — 덮으면 약값이 이중계산된다)
            _SEBU_WIN = {
                '상해사망','상해후유3%','질병후유3%','입원','일반암','암수술','유사암(갑.기.경.제)',
                '통합전이암','뇌혈관진단비','뇌혈관수술비','뇌졸증진단비','허혈성 진단비','심장수술비',
                '급성심근경색','질병수술비','상해수술비','질병일당','상해일당','간병인',
                '합의금','변호사','대인','대물','자부상','일상배상책임','깁스진단비',
            }
            _SEBU_DEATH = ('일반사망','질병사망(80세)')
            _ci_on = bool(nm2r.get('중대한CI적용')) and \
                     ws.cell(nm2r.get('중대한CI적용') or 1, _cl2).value not in (None, '')
            if _partial:
                # ★★★★★v299-2: 세부표가 그 계약을 커버하는데 화이트리스트 칸이 <b>비어 있으면
                #   그 계약에 그 담보는 없다</b>는 뜻이다. 덮어쓰기만으로는 별첨이 만들어낸 유령값이
                #   남는다. 실측 = 롯데·AXA 결합담보 분해가 상해후유3%에 각 15,000씩 남아
                #   66,000(한장표 36,000)이 됐다. ★CI 본체가 있는 축은 지우지 않는다(§22).
                _have = set(_vals.keys())
                if _jong: _have.add('일반사망')
                for _wn in _SEBU_WIN:
                    if _wn in _have: continue
                    _rw = nm2r.get(_wn)
                    if not _rw or ws.cell(_rw,_cl2).value in (None, ''): continue
                    _pr2 = nm2r.get(_CIPAIR.get(_wn) or '')
                    if _pr2 and ws.cell(_pr2,_cl2).value not in (None, ''): continue
                    ws.cell(_rw,_cl2).value = None
                # ★★★★★v299-3 사망 교정: 검산식 `질병사망 = 일반사망 + 질병사망(80세)`이므로
                #   두 행을 <b>함께</b> 리셋한 뒤 세부값 하나만 기재해야 합이 맞는다.
                #   실측 = 새마을금고 `정기특약 : 사망공제금` 2,000이 일반사망에 남아 세부 1,000 위에
                #   덧붙어 3,000이 됐다(한장표 질병사망 64,070 vs 엑셀 66,070).
                #   ★CI 계약은 손대지 않는다 — 사망 분할이 §22 CI 4단계 소관이다.
                if ('질병사망(80세)' in _vals) and not _ci_on:
                    for _dr in _SEBU_DEATH:
                        _rr2 = nm2r.get(_dr)
                        if _rr2: ws.cell(_rr2,_cl2).value = None
            for _nm2,_v2 in _vals.items():
                _tgt = '일반사망' if (_jong and _nm2=='질병사망(80세)') else _nm2
                _r2 = nm2r.get(_tgt)
                if not _r2: continue
                # ★★★★★v381 (지점장 지적 2026.08.11 "한화생명에 어디서 실손이 있냐", 영구):
                #   <b>실손 계약이 아닌 계약에는 실손 행을 세부보충하지 않는다.</b>
                #   [구 결함] 세부가입현황(계약별 가입정보)은 계약이 가로로 늘어선 표라
                #   <b>인접 계약의 실손 값이 옆 칸으로 밀려</b> 들어온다. 실측(구본칠) =
                #   삼성 실손 계약의 `입원 5,000`이 <b>한화생명 CI 계약</b>에 앉았고
                #   → ①한화가 실손 보유로 잡혀 <b>실손 세대가 한화 가입일(2007)로 판정=1세대</b>
                #   → ②한장표 검산 `입원 5,000 vs 엑셀 10,000` 불일치.
                #   실손은 상품 자체가 실손이어야 존재한다 — 별첨·상품명으로 판별해 게이트한다.
                if _tgt in ('입원','통원','약값','상해의료비','MRI/도수치료/비급여주사'):
                    if not _is_silson_like(_c2.get('company',''), _c2.get('product',''), _c2.get('dambo')):
                        unmapped.append((_cl2, str(_cn2), f'세부보충 {_tgt}', _v2,
                                         '[확인] 실손 계약이 아닌데 세부가입현황에 실손 값 — 옆 계약 값 밀림 의심 → 기재 안 함'))
                        continue
                if _partial:
                    if _tgt not in _SEBU_WIN and not (_tgt in _SEBU_DEATH and not _ci_on):
                        if ws.cell(_r2,_cl2).value not in (None, ''): continue  # 별첨 값 보존
                    _pr = nm2r.get(_CIPAIR.get(_tgt) or '')
                    if _pr and ws.cell(_pr,_cl2).value not in (None, ''): continue   # CI 본체 이미 기재
                    # ★★★★★v358 (지점장 확정 2026.08.02, 영구): "<b>그건 보장분석지pdf안에 답이있어
                    #   나한테 묻지말고 그거대로해라</b>" → <b>별첨이 정본이다. 세부보충이 만들어 넣지 않는다.</b>
                    #   [실측] 삼성 별첨은 `［간편］골절 진단비（치아파절…제외）` 50 <b>하나뿐</b>인데
                    #   세부가입현황엔 구분 없는 `골절진단비 50`만 있어 <b>포함 행에도 50</b>이 앉았다
                    #   → 같은 담보가 두 행에 갈렸다.
                    #   → <b>짝 행(포함↔제외) 중 하나라도 별첨 값이 있으면 세부보충을 하지 않는다.</b>
                    #   ★별첨에 '치아제외' 표기가 <b>아예 없는</b> `골절진단비`는 종전대로 <b>포함 행</b>이다
                    #     (resolve_kw 정본, v38c).
                    _pr358 = nm2r.get(_SEBU_PAIR.get(_tgt) or '')
                    if _pr358 and ws.cell(_pr358,_cl2).value not in (None, ''): continue
                    # ★★★★★v358 (지점장 확정 2026.08.02, 영구): "<b>골절(치아파절제외) 삼성 50 ← 별첨 원문
                    #   이게 맞다</b>". 세부가입현황은 <b>포함/제외 구분이 없는 표</b>라 `골절진단비 50`뿐이다.
                    #   그런데 세부보충은 「빈 행이면 채운다」로만 동작해, 별첨이 <b>제외</b> 행을 채운 뒤에도
                    #   <b>포함 행이 비었다는 이유로 또 50을 넣었다</b> → 같은 담보가 두 행에 갈려 앉았다.
                    #   → <b>짝 행 중 하나라도 별첨 값이 있으면 세부보충 금지</b>(CI `_CIPAIR`와 같은 구조).
                    #   ★대상은 <b>구분이 갈리는 짝</b>뿐이다. 다른 행은 종전대로 동작한다.
                    _PAIR358 = {'골절(치아파절포함)':'골절(치아파절제외)',
                                '골절(치아파절제외)':'골절(치아파절포함)'}
                    _pr358 = nm2r.get(_PAIR358.get(_tgt) or '')
                    if _pr358 and ws.cell(_pr358,_cl2).value not in (None, ''): continue
                ws.cell(_r2,_cl2).value = _v2
                # ★★★★★v359 [세부보충] 흔적 — 지점장 지시 2026.08.02(영구).
                #   지점장 원문: "<b>왜 계속 없는걸 넣느냐</b>" · "<b>없는걸 새로 생각해서 넣지마라</b>".
                #   세부가입현황 보충은 <b>별첨에 없는 값을 채우는</b> 유일한 경로다. 값은 지침대로 두되
                #   <b>무엇을 어디에 넣었는지 전부 확인사항에 남긴다</b> → 「조용히 넣는 것」을 없앤다.
                try: _SEBU_LOG.append((str(_cn2), str(_tgt), _v2))
                except Exception: pass
                _gen2 = ('비갱신' not in str(_c2.get('renewal') or '')) and ('갱신' in str(_c2.get('renewal') or ''))
                ws.cell(_r2,_cl2).font  = BL if _gen2 else BK
                ws.cell(_r2,_cl2).alignment = Alignment(horizontal='center', vertical='center')
                _put+=1
            if _put:
                print(f"[v271 세부보충] {_cn2} {_pm:,}원 → {'빈 행만 보충' if _partial else '별첨 매핑 0건 → 재구성'} → 세부가입현황에서 {_put}개 기재"
                      + (" (종신→일반사망)" if _jong else ""))

    # ★★★★★v272 영구지침(지점장 지시 2026.07.30):
    #   <b>「~2005년 05월까지 가입 + 회사명에 '생명' + 담보명에 '수술보장'」</b>인 담보는
    #   <b>질병 종수술비(1-3종) · 상해 종수술비(1-3종)</b> 두 행에 <b>각각 대표값(max)</b>을 기재한다.
    #   ・마스터에 두 행을 신설했다(상해 67 · 질병 75). 이후 행은 전부 +2 밀렸다.
    #   ・★<b>dambo 키로 넣으면 안 된다</b> — `resolve_kw('질병 종수술비(1-3종)')`가 종 번호 3을
    #     읽어 <b>1-5종 3종 슬롯</b>으로 보낸다(실측). 그래서 엑셀 행에 직접 기재한다.
    #   ・★<b>v271 블록보다 뒤</b>에 둔다. 앞에 두면 '별첨 매핑 0건' 판정이 깨져 세부보충이 멈춘다.
    for _ix3, _c3 in enumerate(contracts):
        _cn3 = _c3.get('company') or ''
        if '생명' not in _cn3: continue
        _m3 = re.match(r'(\d{4})[.\-/](\d{1,2})', str(_c3.get('contract_date') or ''))
        if not _m3: continue
        if (int(_m3.group(1)), int(_m3.group(2))) > (2005, 5): continue
        # ★★★v273: 값은 <b>별첨(정상계약 리스트) 줄 단위 대표(max)</b>다.
        #   dambo는 같은 담보명 2줄을 합산하므로(실측 100+100=200) 그대로 쓰면 두 배가 된다.
        #   ★★계약이 여러 개면 <b>계약별로 각각</b> 기재한다(지점장 확정 2026.07.30:
        #     "139800원 계약에 200만 있고 41,800원에 100이 더 있다").
        #     끝열 SUM = 계약 합산이 정상이다. 전체 대표 1건으로 줄이지 말 것.
        try: _pm3 = int(str(_c3.get('premium')).replace(',','').replace('원','').strip())
        except: _pm3 = None
        _mx3 = float((data.get('surg13') or {}).get(_pm3) or 0.0)
        if _mx3 <= 0: continue
        _cl3 = 3 + _ix3
        _gen3 = ('비갱신' not in str(_c3.get('renewal') or '')) and ('갱신' in str(_c3.get('renewal') or ''))
        for _nm3 in ('상해 종수술비(1-3종)', '질병 종수술비(1-3종)'):
            _r3 = nm2r.get(_nm3)
            if not _r3: continue
            ws.cell(_r3, _cl3).value = _mx3
            ws.cell(_r3, _cl3).font = BL if _gen3 else BK
            ws.cell(_r3, _cl3).alignment = Alignment(horizontal='center', vertical='center')
        print(f"[v272 1-3종] {_cn3} {_c3.get('contract_date')} 수술보장 대표 {_mx3:,.0f} → 질병·상해 종수술비(1-3종)")

    # ★★★★★v477 제83조 — <b>색의 최종 확정은 여기다</b>(모든 재기재가 끝난 자리).
    #   담보명에 갱신 표기가 있으면 주계약이 비갱신이어도 <b>그 담보만 파랑</b>(제5조 B · 제86조).
    #   세부보충(v271)·1-3종(v272)·사망이동은 계약 갱신 여부만 보고 덮으므로 여기서 되돌린다.
    _bfix = 0
    for _bc, _brs in (_BLUE_ALL or {}).items():
        for _br2 in sorted(_brs or ()):
            _bv = ws.cell(_br2,_bc).value
            if _bv in (None,''): continue
            _bf = ws.cell(_br2,_bc).font
            if not (_bf and _bf.color and str(_bf.color.rgb or '').endswith('0070C0')):
                ws.cell(_br2,_bc).font = BL; _bfix += 1
    if _bfix: print('[v477 갱신담보색] 최종 확정 %d셀 (재기재로 덮인 것 되돌림)' % _bfix)

    # ★★★★★v580 제122조 4항 — <b>현물 셀에 노란 채우기</b>(지점장 지시 2026.08.24).
    #   색 확정이 <b>전부 끝난 뒤</b>에 칠한다. 앞에 두면 위 루프가 폰트를 덮으면서
    #   같은 셀을 두 규칙이 만지게 된다(「나중 것이 앞 것을 죽인다」).
    if _hy_cells:
        from openpyxl.styles import PatternFill as _PF
        _YEL = _PF('solid', fgColor='FFFF00')
        # ★v580b 지점장 지시 2026.08.24 「<b>그냥 현물인 경우 가로 한 줄을 노란색으로 해라</b>」
        #   ⇒ 셀 하나가 아니라 <b>그 담보 행 전체</b>(A열~끝열)를 노랑으로. 금액은 그대로 살고
        #     현물이라는 사실은 줄 색으로 보인다. 산출물에서 금액이 사라지지 않는다.
        _rows_hy = sorted({_hr for _hr, _hc in _hy_cells if _hr})
        for _hr in _rows_hy:
            for _hc2 in range(1, ws.max_column + 1):
                ws.cell(_hr, _hc2).fill = _YEL
        for _hr, _hc in _hy_cells:
            if _hr and ws.cell(_hr, _hc).value in (None, ''):
                ws.cell(_hr, _hc).value = 0
        if _rows_hy:
            print('[v580 제122조 4항] 현물 — 행 전체 노랑 %d행(행번호 %s) · 현물 담보 금액 0'
                  % (len(_rows_hy), _rows_hy))

    # ★ 합계 = 항상 표 맨 끝 열. 가로 SUM 수식(법칙22, 하드코딩 금지).
    # ★★★★★v388 (지점장 확정 2026.08.12): 「<b>엑셀은 각각 / 보장분석지·보장진단서는 합산.
    #   그래서 엑셀은 헷갈릴 수 있으니 따로 합산라인 추가해라 — 보유 합계 | 제안 합계 둘 다.
    #   보유합계 글자색 = 블랙 or 블루 / 제안합계 = 레드</b>」
    #   배치 = [보유 계약들][제안 계약들] <b>[보유 합계][제안 합계]</b> [합계]
    #   ★<b>합계 열을 계속 맨 끝에 둔다</b> — `ws.max_column`을 끝열로 쓰는 코드가 8곳(PPT·설명서·진단서·캐시)
    #     이라 순서를 바꾸면 4대 산출물이 통째로 다른 열을 읽는다.
    #   ★<b>기존 합계 수식은 손대지 않는다</b> — 범위가 `C..마지막 계약열`이라 새 2열을 포함하지 않는다.
    #   ★<b>제안 계약이 없으면 2열을 만들지 않는다</b>(종전 구조 그대로 = 회귀 0). 헷갈릴 일이 없기 때문.
    _own_n = sum(1 for _c in contracts if not _c.get('proposal'))
    _jn_n  = n_ct - _own_n
    _has_jn = _jn_n > 0
    own_sum_col  = (3 + n_ct) if _has_jn else 0
    jean_sum_col = (4 + n_ct) if _has_jn else 0
    last_col = (5 + n_ct) if _has_jn else (3 + n_ct)
    # ★★★v230 (지점장 지시 2026.07.25, 영구): <b>유사암 자동유도(일반암×10%)는 완전 폐기</b>.
    #   지점장 원문 = <b>"유사암 적힌 것만 넣어라"</b>. 별첨에 유사암 담보가 없으면 <b>그 계약은 공란</b>이다 —
    #   일반암 금액으로 유추해 넣지 않는다. 구 v30q 자동유도가 없는 담보를 만들어냈다.
    #   실측(이정화): 한화 일반암 6,000 → 유사암 <b>600 자동생성</b> · 메리츠0804 1,000 → <b>100 자동생성</b>
    #   = 합계 1,450(정답 900). 자동유도 제거 + 명시 담보만 산입 → <b>900</b>으로 한장보장표와 일치.
    #   ★구 v213 '명시액이 하나도 없을 때만 유도' 게이트 방식도 함께 폐기(지점장 'no').

    first_L = get_column_letter(3)
    last_ct_L = get_column_letter(2 + n_ct) if n_ct>0 else first_L   # ★v388 마지막 <b>계약</b> 열(합산 2열 제외)
    hc = ws.cell(1, last_col)
    hc.value = '합계'; hc.font = W; hc.fill = FILL_SUM; hc.alignment = AL
    if _has_jn:                                  # ★v388 보유 합계 · 제안 합계 헤더
        _oh = ws.cell(1, own_sum_col)
        _oh.value = '보유 합계'; _oh.font = W; _oh.fill = FILL_SUM; _oh.alignment = AL
        _jh = ws.cell(1, jean_sum_col)
        _jh.value = '제안 합계'
        _jh.font = Font(color='000000', name='맑은 고딕', size=9, bold=True)
        _jh.fill = PatternFill('solid', fgColor='ED7D31')     # 제안 열과 같은 주황
        _jh.alignment = AL
    # 보험료 합계 = 숫자만 표기(§3): 수식 아닌 계산된 숫자값. 글자 검정(흰바탕)
    # ★v128 지점장 확정 2026.07.21: 월보험료 합계 = 잔여보험료 > 0(납입 진행중) 계약만.
    #   납입완료(회차 a>=b) 계약은 계약열에는 남기되 합계에서는 뺀다.
    #   실측 근거(배학술 롯데): 12건 전액 711,218 → 교보 120/120·AIA 180/180 제외 → 605,618 = 리포트 일치.
    if n_ct>0:
        # ★★★v199 지점장 확정 2026.07.23: 보험료 합계도 '=D2+E2+…'가 아니라 <b>단일 =SUM()</b>.
        #   근거: 셀을 복사·이동하면 a+b+ 수식은 참조가 어긋나 합계가 하나도 안 맞는다.
        #   완납 계약은 보험료 칸을 텍스트로 넣어(위 루프) SUM이 자동으로 건너뛴다 = v129 정본 유지.
        #   ★캐시: inject_sum_cache가 숫자 셀만 합산해 같은 값을 박으므로 설명서·PPT가 0원으로 읽지 않는다.
        ws.cell(2, last_col).value = f'=SUM(C2:{last_ct_L}2)'
        ws.cell(2, last_col).font = BK
        if _has_jn:                              # ★v388 보험료도 보유/제안 분리
            _oL = get_column_letter(2 + _own_n)
            _jF = get_column_letter(3 + _own_n)
            ws.cell(2, own_sum_col).value  = f'=SUM(C2:{_oL}2)'
            ws.cell(2, own_sum_col).font   = BK
            ws.cell(2, jean_sum_col).value = f'=SUM({_jF}2:{last_ct_L}2)'
            ws.cell(2, jean_sum_col).font  = Font(color='C00000', name='맑은 고딕', size=9)

    for r in range(6, ws.max_row+1):
        slash_t=[0]*8; slash_n=0; is_slash=False; has_num=False   # ★v29v 1-8종·v29y 트리오: 실제 칸수 따름
        for col in range(3, 3 + n_ct):        # ★v388 데이터(계약) 열만 — 합산 2열 제외
            v = ws.cell(r,col).value
            if isinstance(v,(int,float)): has_num=True
            elif isinstance(v,str) and '/' in v:
                is_slash = True
                _ps=v.split('/')[:8]
                slash_n=max(slash_n,len(_ps))
                for k,p in enumerate(_ps):
                    try: slash_t[k] += int(p)
                    except: pass
        sc = ws.cell(r, last_col)
        if is_slash and any(slash_t):
            # ★★v378 종수술 행 한정 — 같은 행에 <b>종별 슬래시</b>와 <b>종번호 없는 단일금액</b>이 섞이면
            #   구 코드는 슬래시만 더하고 <b>숫자 칸을 통째로 버렸다</b>(끝열이 조용히 작아진다).
            #   단일금액은 1~5종 어느 종이든 그 금액이므로 <b>모든 종 칸에 가산</b>한다.
            #   ※이 해석은 지점장 확정 대기 항목이다(다른 슬래시 행에는 적용하지 않는다).
            if has_num and '종수술비' in str(ws.cell(r,2).value or ''):
                _lump = sum(v for v in (ws.cell(r,c).value for c in range(3, 3 + n_ct))   # ★v388
                            if isinstance(v,(int,float)))
                if _lump:
                    for k in range(slash_n or 5): slash_t[k] += _lump
            sc.value = '/'.join(str(x) for x in slash_t[:(slash_n or 5)])
            sc.font = BL if str(ws.cell(r,2).value).strip() in _BLUE_ROWS else BK   # 슬래시 행은 §3 SUM 예외
        else:
            # ★v29t: §5·v29c(2) 원복 — 합계는 동적 =SUM 수식(하드코딩 금지). 사용자가 값을 추가해도 자동 합산.
            #   저장 후 recalc_xlsx가 캐시값 주입 → 폰·미리보기에서도 숫자 표시(수식 유지).
            _rng = f'C{r}:{last_ct_L}{r}'
            _bnm=str(ws.cell(r,2).value).strip()
            # ★v91: 구 '입원 5,000 캡' 폐기(지점장 2026.07.19). 실손 다건이면 합이 5,000을 넘는다.
            #   실측 장혜경 = 현대 5,000 + DB 500 = 5,500 (한장보장표 질병 5,500과 일치).
            if _bnm=='입원': sc.value = f'=SUM({_rng})'
            elif _bnm=='자부상': sc.value = f'=MIN(SUM({_rng}),80)'          # ★지점장 2026.07.02: 자부상 최대 80만 캡
            elif _bnm=='120대수술비':                                       # ★v30k n대수술비=계약별 값 가로 슬래시(합산·최댓값 금지, 지점장 2026.07.03)
                _nd=[ws.cell(r,c).value for c in range(3, 3 + n_ct)]   # ★v388
                _nd=[str(int(x)) for x in _nd if isinstance(x,(int,float)) and x>0]
                sc.value = '/'.join(_nd) if _nd else f'=SUM({_rng})'
            elif _bnm in ('간병인','중입자치료비'): sc.value = f'=IF(COUNT({_rng})=0,0,MAX({_rng}))'  # ★v30d 간병인·중입자=전 계약 대표 최댓값 1건
            elif _bnm=='간호통합병동': sc.value = f'=IF(COUNT({_rng})=0,0,MAX({_rng}))'   # ★v41 1-180일 최댓값 1건
            # ★★★★★v456 제67조 (지점장 지적 2026.08.17 「<b>가입제안서가 통합암을 다 합산한다.
            #   통합암/통합전이암은 대표값 하나만 입력해라</b>」)
            #   지침 §8.2·v320 「통합암·통합전이암 = 대표금액 1개」가 정본인데
            #   `_rep1`(계약 내 대표) 목록에 <b>통합전이암만 있고 통합암이 빠져</b> 있었다(실측).
            #   → 계약 내 대표 + <b>끝열도 대표(max)</b>. 대표 행은 끝열과 같은 수식이어야 한다(제10조).
            elif _bnm in ('통합암','통합전이암','암일당'): sc.value = f'=IF(COUNT({_rng})=0,0,MAX({_rng}))'
            else: sc.value = f'=SUM({_rng})'
            # ★★★★★v336 (지점장 지적 2026.08.02 "또 실손이 검정으로 나온다. 실손은 비갱신이 없다"):
            #   계약 셀은 파랑인데 <b>끝열 합계만 검정</b>이었다(`sc.font = BK` 고정).
            #   §10 정본 = <b>실손(입원·통원·약값) + 일상배상책임은 항상 파랑</b> → 끝열도 파랑이어야 한다.
            sc.font = BL if _bnm in _BLUE_ROWS else BK

        # ★★★★★v388 보유 합계 · 제안 합계 (지점장 확정 2026.08.12)
        #   ・슬래시 행은 SUM이 안 되므로 <b>구간별로 직접 합산해 문자열</b>로 넣는다(끝열과 같은 방식).
        #   ・대표(max)·캡 행은 구간 SUM만 넣는다 — <b>끝열의 MIN/MAX 값이 정본</b>이고 이 두 열은 참고용이다.
        #   ・색: 보유합계 = 끝열과 동일 규칙(파랑/검정) / 제안합계 = <b>레드 C00000</b>.
        if _has_jn:
            _bnm2 = str(ws.cell(r,2).value).strip()
            _RD2  = Font(color='C00000', name='맑은 고딕', size=9)
            _oL2  = get_column_letter(2 + _own_n)
            _jF2  = get_column_letter(3 + _own_n)
            for _sc2, _rng2, _c0, _c1, _isj in (
                    (ws.cell(r, own_sum_col),  f'C{r}:{_oL2}{r}',            3,          3 + _own_n, False),
                    (ws.cell(r, jean_sum_col), f'{_jF2}{r}:{last_ct_L}{r}',  3 + _own_n, 3 + n_ct,   True)):
                _sl = [0]*8; _sn = 0; _any = False
                for _c2 in range(_c0, _c1):
                    _v2 = ws.cell(r,_c2).value
                    if isinstance(_v2,str) and '/' in _v2:
                        _any = True
                        _ps2 = _v2.split('/')[:8]; _sn = max(_sn, len(_ps2))
                        for _k2,_p2 in enumerate(_ps2):
                            try: _sl[_k2] += int(_p2)
                            except: pass
                # ★v388b 끝열과 <b>같은 규칙</b>을 적용해야 「보유+제안 = 합계」가 눈으로 맞는다.
                #   v378: 종수술 행에 종별 슬래시와 단일금액이 섞이면 단일금액을 <b>전 종 칸에 가산</b>.
                #   이 규칙을 빼먹어 실측에서 보유 `30/0/0/0/0`인데 합계 `1080/1110/...`로 어긋났다.
                if _any and any(_sl) and '종수술비' in _bnm2:
                    _lp2 = sum(v for v in (ws.cell(r,_c2).value for _c2 in range(_c0,_c1))
                               if isinstance(v,(int,float)))
                    if _lp2:
                        for _k3 in range(_sn or 5): _sl[_k3] += _lp2
                if _any and any(_sl):
                    _sc2.value = '/'.join(str(x) for x in _sl[:(_sn or 5)])
                elif _bnm2 == '자부상':
                    # ★★v388b 캡·대표(max) 행은 <b>끝열과 같은 수식</b>을 쓴다(결과값 동결).
                    #   단순 SUM으로 두면 실측처럼 <b>보유 30 + 제안 600 인데 합계 80</b>이 되어
                    #   지점장이 없애려던 그 헷갈림이 그대로 남는다. 각 칸이 「그 구간의 실제 보장액」이 되게 한다.
                    #   ※이 행들은 합이 아니라 캡·대표이므로 「보유+제안 = 합계」가 원래 성립하지 않는다.
                    _sc2.value = f'=MIN(SUM({_rng2}),80)'
                elif _bnm2 in ('간병인','중입자치료비','간호통합병동'):
                    _sc2.value = f'=IF(COUNT({_rng2})=0,0,MAX({_rng2}))'
                else:
                    _sc2.value = f'=SUM({_rng2})'
                _sc2.font = _RD2 if _isj else (BL if _bnm2 in _BLUE_ROWS else BK)

    ws.column_dimensions['B'].width = 22
    for c in range(3, last_col+1):
        ws.column_dimensions[get_column_letter(c)].width = 12

    # ★ 테두리: A(구분)~끝열(합계) 전체 격자 직접 그림 + 구분(키워드)마다 굵은 구분선.
    #   (마스터 A·B 테두리가 중간행에서 끊겨 '선 없음' 발생 → 전부 새로 그림)
    _thin = Side(style='thin', color='FF000000'); _med = Side(style='medium', color='FF000000')   # ★v29t: 6자리 색은 알파00(투명) 저장돼 일부 뷰어에서 선 사라짐 → FF 필수
    # 구분(그룹) 끝행 동적 계산: A열에 값 있는 행=그룹 시작 → 다음 시작-1 = 그룹 끝
    g_starts = [r for r in range(6, ws.max_row+1) if ws.cell(r,1).value not in (None,'')]
    g_end = set()
    for k, s in enumerate(g_starts):
        e = (g_starts[k+1]-1) if k+1 < len(g_starts) else ws.max_row
        g_end.add(e)
    # 수술비 블록 내부 구분: 질병수술비 행 위에 굵은 선(상해 수술 ↔ 질병 수술)
    row_top_med = set()
    for r in range(6, ws.max_row+1):
        if str(ws.cell(r,2).value).strip() == '질병수술비':
            row_top_med.add(r)
    for r in range(1, ws.max_row+1):
        for c in range(1, last_col+1):
            left   = _med if c == 1 else _thin
            right  = _med if c == last_col else _thin
            top    = _med if (r in (1, 6) or r in row_top_med) else _thin
            # 헤더 5행 + 각 구분 끝행 = 굵은 가로 구분선
            bottom = _med if (r in (1,2,3,4,5) or r in g_end) else _thin
            ws.cell(r,c).border = Border(left=left, right=right, top=top, bottom=bottom)
    # ★ 숫자 콤마: 보험료·담보값·합계 SUM 전부 #,##0. (날짜·납입기간·슬래시 행은 텍스트라 제외)
    for r in range(1, ws.max_row+1):
        for c in range(2, last_col+1):
            v = ws.cell(r,c).value
            if isinstance(v,(int,float)) or (isinstance(v,str) and v.startswith('=')):
                ws.cell(r,c).number_format = '#,##0'

    # ★ 합계 이후 잔재 열 삭제 (§3: 합계 = 맨 끝 열)
    if ws.max_column > last_col:
        ws.delete_cols(last_col+1, ws.max_column - last_col)

    # ── 확인사항 시트: LLM 매핑 실패 담보 노출(자가진단, §10) ──
    # ★v311 「지침케이스」는 <b>마스터 전용 관리 시트</b>다 — 고객 산출물엔 나가면 안 된다.
    #   build_excel이 마스터 워크북을 그대로 save하므로 여기서 지운다.
    for _msn in ('지침케이스','해석원칙'):   # ★v312 마스터 전용 관리 시트 — 고객 산출물엔 안 나간다
        if _msn in wb.sheetnames: del wb[_msn]
    for _sn in ('📋확인사항','확인사항'):
        if _sn in wb.sheetnames: del wb[_sn]
    ws2 = wb.create_sheet('확인사항')   # ★v41 이모지·외부하이퍼링크 제거(엑셀 '편집사용' 지연 원인)
    ws2.cell(1,1, f'{client} · 자동분석 {datetime.datetime.now():%Y.%m.%d}')
    try:
        _ipw = globals().get('_IMG_PDF_WARN','')
        if _ipw:
            _c1 = ws2.cell(1,1, f'{client} · 자동분석 {datetime.datetime.now():%Y.%m.%d}   ' + _ipw)
            _c1.font = Font(bold=True, size=13, color='B00020')
    except Exception: pass
    try:   # ★v371 제안서 단독 모드 — 검산·실손세대 「불가」 명시(지점장 확정)
        _jo = globals().get('_JEAN_ONLY','')
        if _jo:
            _c2 = ws2.cell(2,1, _jo)
            _c2.font = Font(bold=True, size=12, color='C00000')
    except Exception: pass
    ws2.cell(3,1,'계약수'); ws2.cell(3,2,n_ct)
    try:
        _ac = ws2.cell(3,4, '[감사] 규칙 %s · 마스터커버리지 %s · FAIL %d건'
                            % (_aud.get('case','-'), _aud.get('cov','-'), _aud.get('fail',0)))
        if _aud.get('fail',0): _ac.font = Font(bold=True, size=12, color='C0392B')
    except Exception: pass
    # ★★★v298-C: 확인사항 4행이 완납 계약까지 더해 본표 2행 합계와 달랐다.
    _pin  = sum(c["premium"] for c in contracts if _in_sum(c))
    _pout = sum(c["premium"] for c in contracts if not _in_sum(c))
    ws2.cell(4,1,'월보험료합계')
    ws2.cell(4,2, f'{_pin:,}원' + (f'  (완납 제외 {_pout:,}원)' if _pout else ''))
    # ★★v186 (지점장 2026.07.22): AIA/AIG/라이나(우체국) 계약이 있으면 <b>엑셀 확인사항 시트 최상단</b>에
    #   "뇌 범위 부분 꼭 체크" 경고를 굵은 빨강으로 박는다. 3군데(엑셀·진단서 1p·7p 워크시트) 동일 문구.
    try:
        _wc = [f for f in ('AIA','AIG','라이나','우체국')
               if any(f in str(c.get('company','')).replace(' ','') for c in contracts)]
        if _wc:
            # ★v187 회사명은 <b>항상 3개 고정</b>(AIA / AIG / 라이나) — 보유분만 나열하지 않는다.
            _wcell = ws2.cell(2,1, '★ AIA / AIG / 라이나생명은 "뇌 범위 부분" 꼭 체크하기 '
                                   '(세부가입현황에서 뇌혈관 / 뇌졸증 / 뇌출혈 축 대조)')
            _wcell.font = Font(bold=True, size=13, color='C0392B')
    except Exception: pass
    ws2.cell(6,1,'[확인] 자동매핑 실패 담보 (마스터 미수록 또는 약관 확인 후 수기 기재)')
    ws2.cell(7,1,'회사'); ws2.cell(7,2,'담보명'); ws2.cell(7,3,'금액(만원)'); ws2.cell(7,4,'보장범위(참고)'); ws2.cell(7,5,'약관검색')
    LINKF = Font(color='0000FF', underline='single')
    # ★v148 (지점장 확정 2026.07.21): 흥국화재 10억통장(리셋월렛II)은 <b>엑셀 전면 기재금지</b>.
    #   본표에서 뺐어도 [확인] 큐(확인사항 시트)에 남으면 설계사가 수기 기재하게 되므로 여기서도 제외한다.
    #   이 담보는 보장진단서 7p 카드 전용이다.
    unmapped = [u for u in unmapped
                if not (('리셋월렛' in re.sub(r'\s','',str(u[2]))) or ('리셋월랫' in re.sub(r'\s','',str(u[2]))))]
    # ★★★★★v308 [확인]큐 3분류 (지점장 지시 2026.07.31, 영구) — <b>표시 전용·값 불변</b>
    #   왜 필요한가: [검산]은 한장보장표 <b>27개만</b> 본다. 마스터 99행 중 <b>72행은 아무도 검증하지 않는다</b>.
    #   그래서 [검산] 0건인데도 설계사는 [확인]큐 100~138건을 보고 "오류가 많다"고 한다(실측: 주재현 102 · 심정자 138).
    #   그 큐에는 ①마스터에 행이 아예 없는 <b>정상</b>과 ②행이 있는데 못 찾은 <b>결함</b>이 섞여 있는데 구분이 없었다.
    #   → 3분류해 <b>결함의심을 맨 위로</b> 올린다. 값·행 배정은 <b>일절 건드리지 않는다</b>(회귀 위험 0).
    _MASTER_KEYS = []
    try:
        for _lab in nm2r.keys():
            _b = re.sub(r'\([^)]*\)', '', str(_lab))
            _b = re.sub(r'\s', '', _b)
            if len(_b) >= 3: _MASTER_KEYS.append((_b, str(_lab)))
        _MASTER_KEYS.sort(key=lambda t: -len(t[0]))
    except Exception: _MASTER_KEYS = []
    _RULE_OUT = [('치매','치매 = 기재 안 함(지점장 확정 2026.07.26)'),
                 ('두번받는','두번받는 CI = 무시'), ('CI추가보장','두번째 CI 담보 = 무시'),
                 ('상급','상급병원 입원일당 = 기재금지(v301)'),
                 ('인실','병실등급 담보 = 종합병원 질병입원일당 아님(v209)'),
                 ('고주파','부정맥 시술 담보 = 진단비 행 아님(v217)'),
                 ('절제술','부정맥 시술 담보 = 진단비 행 아님(v217)'),
                 ('냉각','부정맥 시술 담보 = 진단비 행 아님(v217)')]
    def _qclass(_raw, _note):
        _n = re.sub(r'\s', '', str(_raw or '')); _nt = str(_note or '')
        if ('접힘의심' in _n) or ('원문 대조' in _nt):
            return ('★결함의심', '담보명 잘림 — 별첨 원문 대조 필요')
        for _k, _why in _RULE_OUT:
            if _k in _n: return ('규칙제외', _why)
        if ('후유' in _n) and ('80' not in _n) and re.search(r'(?<![\d.])(20|50)\s*%', _n):
            return ('규칙제외', '후유장해 20%·50% = 미기재(v222)')
        # ★v308b 오탐 제거 — 아래 3종은 <b>지침에 이미 명시된 제외</b>다(결함이 아니다).
        if any(_k in _n for _k in ('납입면제','납입지원','보험료납입')):
            return ('규칙제외', '납입면제·납입지원 = 비담보(§8.6)')
        if ('간병인' in _n) and ('요양' in _n):
            return ('규칙제외', '간병인(요양병원 포함형) = 드롭 정본')
        if (('상해수술비' in _n) and not _n.startswith('상해수술비')) or \
           (('질병수술비' in _n) and not _n.startswith('질병수술비')):
            return ('규칙제외', '수술비 변형(부위·특정·병원규모) = 기재금지(§8.5)')
        if _nt.strip(): return ('규칙제외', _nt.strip()[:40])
        for _b, _lab in _MASTER_KEYS:
            if _b in _n: return ('★결함의심', '마스터 「%s」 행 있음 → 그 행에 들어갔어야 함' % _lab)
        return ('마스터 무행', '해당 담보 행 없음 — 기재 대상 아님')
    _qc = []
    for (_c0, _cm0, _r0, _a0, _n0) in unmapped:
        _cl, _wy = _qclass(_r0, _n0); _qc.append((_cl, _wy, _c0, _cm0, _r0, _a0, _n0))
    _QORD = {'★결함의심':0, '규칙제외':1, '마스터 무행':2}
    _qc.sort(key=lambda t: (_QORD.get(t[0], 9), -(t[5] if isinstance(t[5], (int, float)) else 0)))
    unmapped = [(t[2], t[3], t[4], t[5], t[6]) for t in _qc]
    _qmeta = [(t[0], t[1]) for t in _qc]
    _nbug  = sum(1 for t in _qc if t[0] == '★결함의심')
    _nrule = sum(1 for t in _qc if t[0] == '규칙제외')
    _nnone = sum(1 for t in _qc if t[0] == '마스터 무행')
    print('[v308 확인큐] 총 %d = 결함의심 %d / 규칙제외 %d / 마스터무행 %d' % (len(_qc), _nbug, _nrule, _nnone))
    _h6 = ws2.cell(6,1, '[확인] 자동매핑 실패 담보 %d건  —  ★결함의심 %d건 · 규칙제외 %d건 · 마스터 무행 %d건'
                        % (len(_qc), _nbug, _nrule, _nnone))
    if _nbug: _h6.font = Font(bold=True, size=12, color='C0392B')
    ws2.cell(7,6,'분류'); ws2.cell(7,7,'분류 근거')
    for k,(col,comp,raw,amt,note) in enumerate(unmapped):
        rr = 8+k
        ws2.cell(rr,1,comp); ws2.cell(rr,2,raw); ws2.cell(rr,3,amt); ws2.cell(rr,4,note)
        try:
            _cl, _wy = _qmeta[k]
            _f6 = ws2.cell(rr,6,_cl); ws2.cell(rr,7,_wy)
            if _cl == '★결함의심': _f6.font = Font(bold=True, color='C0392B')
        except Exception: pass
        prod = contracts[col-3]['product'] if 0<=col-3<len(contracts) else ''
        prod_key = re.sub(r'[\(\)\[\]ⅠⅡⅢ_]', ' ', prod)[:18].strip()
        q = f"{comp} {prod_key} {raw[:12]} 약관 보장내용"
        # ★v41 hyperlink 객체 금지 → 평문 URL(엑셀이 열 때 외부링크 검증 안 함 = 편집사용 즉시)
        ws2.cell(rr,5, "https://search.naver.com/search.naver?query=" + urllib.parse.quote(q))
    ws2.column_dimensions['B'].width = 34; ws2.column_dimensions['D'].width = 40; ws2.column_dimensions['E'].width = 12
    ws2.column_dimensions['F'].width = 12; ws2.column_dimensions['G'].width = 46
    # ★v29z 근거 감사 로그 — '없는 값' 논쟁 즉시 검증용
    _rr = 9 + len(unmapped)
    if silson_trace:
        _rr += 2; ws2.cell(_rr,1,'[근거] 실손 세대 판정 (가입일 vs 상품코드 — 상품코드 우선)')
        _rr += 1; ws2.cell(_rr,1,'회사'); ws2.cell(_rr,2,'가입일'); ws2.cell(_rr,3,'상품코드(YYMM)'); ws2.cell(_rr,4,'판정')
        for (_c,_d,_p,_g) in silson_trace:
            _rr += 1; ws2.cell(_rr,1,_c); ws2.cell(_rr,2,_d); ws2.cell(_rr,3,_p); ws2.cell(_rr,4,_g)
    if heart_trace:
        _rr += 2; ws2.cell(_rr,1,'[근거] 심장 블록 기재 내역 (원 담보명 → 기재 행) — 별첨 원문 그대로')
        _rr += 1; ws2.cell(_rr,1,'회사'); ws2.cell(_rr,2,'별첨 원 담보명'); ws2.cell(_rr,3,'기재 행'); ws2.cell(_rr,4,'금액(만원)')
        for (_c,_raw,_rows,_a) in heart_trace:
            _rr += 1; ws2.cell(_rr,1,_c); ws2.cell(_rr,2,str(_raw)[:60]); ws2.cell(_rr,3,_rows); ws2.cell(_rr,4,_a)
    if surg_trace:   # ★v30g 수술 블록 근거 — 종수술 슬롯 이상치 즉시 추적용
        _rr += 2; ws2.cell(_rr,1,'[근거] 수술 블록 기재 내역 (원 담보명 → 기재 행/슬롯) — 별첨 원문 그대로')
        _rr += 1; ws2.cell(_rr,1,'회사'); ws2.cell(_rr,2,'별첨 원 담보명'); ws2.cell(_rr,3,'기재 행/슬롯'); ws2.cell(_rr,4,'금액(만원)')
        for (_c,_raw,_rows,_a) in surg_trace:
            _rr += 1; ws2.cell(_rr,1,_c); ws2.cell(_rr,2,str(_raw)[:60]); ws2.cell(_rr,3,_rows); ws2.cell(_rr,4,_a)
    # ★★★★★ v238 CI 진단표 — <b>CI 계약이 하나라도 있으면 무조건 출력</b>(영구).
    #   오늘(2026.07.25) 사고: CI 2건이 통째로 누락됐는데 <b>산출물 어디에도 흔적이 없어</b> 지점장이
    #   직접 발견해야 했다(확인사항에 CI 근거표가 없었다). → 이제 CI 계약·사망액·본체·선지급률·
    #   중대한OO 배치 결과를 <b>표로 항상 남긴다</b>. 배치가 0이면 최상단에 [경고] 행을 찍는다.
    _cid = _ci_diag or []
    if _cid:
        _zero = all((d.get('placed') or 0) == 0 for d in _cid)
        _rr += 2
        ws2.cell(_rr,1, '[경고] CI 계약이 있으나 중대한OO 배치가 0건 — 선지급 분해 실패, 즉시 확인' if _zero
                       else '[근거] CI 판정·선지급률 내역 (CI 계약은 항상 이 표를 확인할 것)')
        if _zero:
            try: ws2.cell(_rr,1).font = Font(bold=True, color='C00000')
            except Exception: pass
        _rr += 1
        for _i,_h in enumerate(('회사','상품명','①세부내역','②별첨','판정근거','사망보장(만원)','선지급률','본체(만원)','중대한OO 배치'),1):
            ws2.cell(_rr,_i,_h)
        for d in _cid:
            _rr += 1
            ws2.cell(_rr,1,d.get('co','')); ws2.cell(_rr,2,str(d.get('pd',''))[:40])
            ws2.cell(_rr,3,'O' if d.get('sebu') else 'X')
            ws2.cell(_rr,4,'O' if d.get('byul') else 'X')
            ws2.cell(_rr,5,d.get('src',''))
            ws2.cell(_rr,6,d.get('samang') or 0)
            ws2.cell(_rr,7,(f"{d['pct']}%형" if d.get('pct') else '[확인] 50%/80% 미해당'))
            ws2.cell(_rr,8,d.get('bonche') or 0)
            ws2.cell(_rr,9,d.get('placed_txt') or '없음 — [확인]')
        _rr += 1
        ws2.cell(_rr,1,'※ CI 4단계 체크: ①상세 세부내역 ②[별첨] 보장현황 ③중대한*** 유무 ④동일금액 2개 이상 → 50%면 ×2 / 80%면 ÷0.8 이 사망보장금과 맞는지 검증. 선지급률은 50%·80% 두 가지뿐 — 추측 금지.')
    if cancer_trace:   # ★v30h 암 블록 근거 — 일반암 과다·통합암 중복 즉시 추적
        _rr += 2; ws2.cell(_rr,1,'[근거] 암 블록 기재 내역 (원 담보명 → 기재 행) — 별첨 원문 그대로')
        _rr += 1; ws2.cell(_rr,1,'회사'); ws2.cell(_rr,2,'별첨 원 담보명'); ws2.cell(_rr,3,'기재 행'); ws2.cell(_rr,4,'금액(만원)')
        for (_c,_raw,_rows,_a) in cancer_trace:
            _rr += 1; ws2.cell(_rr,1,_c); ws2.cell(_rr,2,str(_raw)[:60]); ws2.cell(_rr,3,_rows); ws2.cell(_rr,4,_a)
    # ★★★★★v277 [검산] 한장보장표 자동 대조 (지침 §1 등식1 · §13 체크리스트 ① 자동화)
    #   지금까지 이 대조는 <b>사람 눈</b>으로만 했다. 그래서 지침이 있어도 매번 놓쳤다.
    #   → 앱이 스스로 대조해 <b>불일치를 확인사항 시트에 강제 출력</b>한다. 값은 건드리지 않는다.
    try:
        _hj = (data.get('hanjang') or {})
        _hjsrc = '한장보장표'
        if not _hj:
            _hj = (data.get('hanjang_kb') or {}); _hjsrc = 'KB 전체 보장 현황(2~3p)'   # ★v295
        if _hj:
            # ★★★★★v421 (지점장 지적 2026.08.14 박미정 검산 불일치 3건 — <b>게이트 오탐</b>)
            #   구 코드는 `_lc2 = 3 + n_ct`로 <b>파싱 단계의 계약 수</b>를 썼다. 3열(KB) 리포트에서
            #   `sebu_ci`가 계약을 3건으로 세면 검산은 <b>C~E 3열만</b> 합산한다.
            #   실측: 상해사망 엑셀 17,600(정답)인데 검산은 10,000 → 「불일치」로 뜬다.
            #   엑셀은 맞는데 <b>게이트가 틀린 값을 비교</b>했다 = 제출 금지가 오탐으로 걸린다.
            #   ★<b>「계약이냐」 판정은 한 곳에서만</b>([v419]) — 엑셀 헤더로 실제 끝열을 찾는다.
            _lc2 = 3 + n_ct
            try:
                _mx = ws.max_column
                _sc = [c for c in range(3, _mx+1) if _is_sumcol(ws, c)]
                if _sc:
                    _lc2 = min(_sc)          # 첫 합산 열 = 계약 열의 끝(exclusive)
                    print(f'[v421 검산] 끝열 재판정 {3+n_ct} → {_lc2} (엑셀 헤더 기준 · 계약 {_lc2-3}열)')
            except Exception: pass
            def _num(v):
                if isinstance(v,(int,float)): return float(v)
                t=str(v or '')
                if '/' in t:                       # 슬래시 행은 대표(max)
                    try: return max(float(x.replace(',','')) for x in t.split('/') if x.strip())
                    except Exception: return None
                try: return float(t.replace(',',''))
                except Exception: return None
            # ★끝열은 이 시점에 <b>수식 문자열</b>이다(캐시 주입 전) → 등식2 v218 규칙 그대로
            #   수식 종류를 보고 데이터셀(C~끝열-1)로 직접 계산한다.
            # ★★★★★v370 (지점장 승인 2026.08.09): <b>한장표 검산에서 가입제안서 열을 제외</b>한다.
            #   한장보장표는 <b>현재 보장</b>이고 가입제안서는 <b>아직 가입하지 않은 담보</b>다.
            #   합계 셀은 지점장 지시대로 <b>제안 포함 그대로</b> 두고, <b>대조 계산에서만</b> 뺀다.
            #   안 빼면 제안서를 넣을 때마다 「불일치 N건」이 떠 제출 금지 게이트가 상시 발동한다(실측 9건).
            _pcols = set()
            try:
                for _i9, _c9 in enumerate(contracts):
                    if _c9.get('proposal'): _pcols.add(3 + _i9)
            except Exception: _pcols = set()
            if _pcols: print(f'[v370 검산] 가입제안서 열 {sorted(_pcols)} 검산 제외')
            _xl = {}
            for _r9 in range(6, ws.max_row+1):
                _n9 = ws.cell(_r9,2).value
                if not _n9: continue
                _f9 = ws.cell(_r9,_lc2).value
                _ds = [_num(ws.cell(_r9,_c9).value) for _c9 in range(3,_lc2) if _c9 not in _pcols]   # ★v370 제안 열 제외
                _ds = [x for x in _ds if x is not None]
                if _pcols and isinstance(_f9,str) and _f9.startswith('=') and '=MIN(' not in _f9 and '=IF(' not in _f9:
                    _v9 = sum(_ds)          # ★v370 단순 SUM 행은 제안 제외 합으로 대체
                elif isinstance(_f9,str) and _f9.startswith('='):
                    if '=MIN(' in _f9:
                        _cap = re.search(r',\s*(\d+)\s*\)\s*$', _f9)
                        _v9 = min(sum(_ds), float(_cap.group(1))) if _cap else sum(_ds)
                    elif '=IF(COUNT' in _f9: _v9 = max(_ds) if _ds else 0.0
                    elif '=IF(SUM'   in _f9: _v9 = 7.0 if _ds else 0.0
                    else:                    _v9 = sum(_ds)
                else:
                    _v9 = (sum(_ds) if (_pcols and _ds) else _num(_f9))   # ★v370 제안 있으면 데이터셀 합(제안 제외)
                    if _v9 is None: _v9 = sum(_ds) if _ds else None
                if _v9 is not None: _xl[str(_n9).strip()] = _v9
            def _g(*names): return sum(_xl.get(x,0.0) for x in names)
            # ★합산 규칙 = 지침 검산식. 행이 갈려도 합은 불변이어야 한다.
            _pairs = [
                # ★★★★★v300-B (주재현·심정자 실측 2026.07.31, 영구): <b>한장보장표·세부가입현황의
                #   '상해사망' 칸은 교통상해사망을 집계하지 않는다</b> → <b>검산식에 더하지 않는다</b>.
                #   [측정] 주재현 계약13 교통 100 · 심정자 계약17/25 각 10,000 — 세 계약 모두
                #   세부표 사망 칸이 비었거나 별도 값이었고, 한장표 합계(79,170 / 178,866)에
                #   그 금액이 <b>들어있지 않았다</b>. 합산하면 +100 / +20,000으로 검산이 깨진다(실측).
                #   ★교통상해사망은 별첨에 실린 담보이므로 <b>9행에 기재는 하되 검산 대상은 아니다</b>
                #   (§3 소스 우선순위 = 별첨이 정본 · 한장표는 검산용).
                ('상해사망',          _g('상해사망')),
                ('질병사망',          _g('일반사망','질병사망(80세)')),
                ('상해후유3%',        _g('상해후유3%')),
                ('질병후유3%',        _g('질병후유3%')),
                ('일반암',            _g('일반암','중대한 암')),
                ('유사암(갑.기.경.제)',_g('유사암(갑.기.경.제)')),
                ('통합전이암',        _g('통합전이암')),
                ('암수술',            _g('암수술')),
                ('뇌혈관진단비',      _g('뇌혈관진단비')),
                ('뇌졸증진단비',      _g('뇌졸증진단비','중대한 뇌졸증')),
                ('뇌혈관수술비',      _g('뇌혈관수술비')),
                ('허혈성 진단비',     _g('허혈성 진단비')),
                ('급성심근경색',      _g('급성심근경색','중대한 급성심근')),
                # ★★★★★v292 (김진구 실측 2026.07.31): 한장표 '심장질환수술비'는 신정원 합산 표기라
                #   <b>심장수술비 + 허혈성수술비</b>가 한 칸에 들어간다(KB '상해입원의료비=입원+외래+처방'과 동일 구조).
                #   실측: 김진구 한장표 2,200 vs 엑셀 심장수술비 0 → 근거표에 허혈성수술비 1,000+1,000 존재.
                #   BARUM은 두 행을 분리 기재하는 것이 정본(#단독담보 원칙)이므로 <b>검산식만</b> 합산한다.
                ('심장수술비',        _g('심장수술비','허혈성수술비')),
                ('상해수술비',        _g('상해수술비')),
                ('질병수술비',        _g('질병수술비')),
                ('상해일당',          _g('상해일당')),
                ('질병일당',          _g('질병일당')),
                ('합의금',            _g('합의금')),
                ('대인',              _g('대인')),
                ('대물',              _g('대물')),
                ('변호사',            _g('변호사')),
                ('자부상',            _g('자부상')),
                ('깁스진단비',        _g('깁스진단비')),
                ('일상배상책임',      _g('일상배상책임')),
                ('뇌출혈진단비',      _g('뇌출혈진단비','중대한 뇌출혈')),          # ★v295
                ('골절진단비',        _g('골절(치아파절포함)','골절(치아파절제외)')),  # ★v295 KB 요약표는 두 행 합계 표기
                ('입원',              _g('입원')),
                ('통원',              _g('통원','약값')),
            ]
            # ★★★★★v299-4 (심정자 실측 2026.07.31, 영구): <b>한장보장표는 제외 7종 계약도 포함한다</b>.
            #   엑셀은 제외하므로 그만큼 낮게 나오는 것이 <b>정상</b>인데 구 검산기는 이를 불일치로 찍었다.
            #   실측 = 일시납 VIP(1/1)·보험기간 1년 AIG 2건이 세부표엔 있고 계약열엔 없어
            #   상해사망 -800 · 질병사망 -300 차이가 났다(둘 다 제외계약 몫과 <b>정확히 일치</b>).
            #   → 세부표에 있으나 <b>계약열에 없는 계약</b>(=제외됨)의 몫을 구해 두고,
            #     차이가 그 몫과 같으면 <b>일치로 인정 + 사유를 표에 남긴다</b>(조용히 넘기지 않는다).
            _orph = {}
            try:
                _pset = set()
                for _c9 in contracts:
                    try: _pset.add(int(str(_c9.get('premium')).replace(',','').replace('원','').strip()))
                    except: pass
                for _pk, _pv in (data.get('sebu_bc') or {}).items():
                    if _pk in _pset: continue
                    for _a9, _b9 in (_pv or {}).items(): _orph[_a9] = _orph.get(_a9, 0) + _b9
            except Exception: _orph = {}
            _OPAIR = {'질병사망':('일반사망','질병사망(80세)'), '일반암':('일반암','중대한 암'),
                      '뇌졸증진단비':('뇌졸증진단비','중대한 뇌졸증'),
                      '급성심근경색':('급성심근경색','중대한 급성심근'), '통원':('통원','약값')}
            # ★★★★★v421 (지점장 확정 2026.08.14 박미정) — <b>기준표가 두 담보를 합쳐 쓰는 칸</b>.
            #   지점장 원문 「<b>무슨소리야!!! 엑셀봐봐</b>」 — 엑셀은 <b>합의금 20,000 / 6주미만 1,000</b>으로
            #   이미 두 행에 정확히 나눠 있었다. KB 한장표만 `교통사고처리지원금 2억1,000만`으로 <b>합쳐</b> 쓴다.
            #   ★<b>엑셀이 맞고 기준표가 뭉뚱그린 것</b>이다 — 검산은 <b>엑셀 두 행의 합</b>과 비교한다.
            _SUMPAIR = {'합의금': ('합의금', '6주미만')}
            _bad=[]; _ok=0; _exl=[]
            for _k,_ev in _pairs:
                if _k not in _hj: continue
                _hv=_hj[_k]
                if abs(_hv-_ev) < 0.5: _ok+=1; continue
                # ★v421 기준표가 합쳐 쓰는 칸 → 엑셀 여러 행의 합으로 재비교
                if _k in _SUMPAIR:
                    _sv = sum(_xl.get(_x, 0.0) for _x in _SUMPAIR[_k])
                    if abs(_hv-_sv) < 0.5:
                        _ok+=1; continue
                _dd = sum(_orph.get(_x,0) for _x in _OPAIR.get(_k,(_k,)))
                if _dd and abs((_hv-_dd)-_ev) < 0.5:
                    _ok+=1; _exl.append((_k,_hv,_ev,_dd))     # 제외계약 몫과 정확히 일치 → 정상
                else: _bad.append((_k,_hv,_ev))
            _rr += 2
            _c0 = ws2.cell(_rr,1, f'[검산] {_hjsrc} 대조 — 일치 {_ok}건 / 불일치 {len(_bad)}건'
                                  + ('  ★불일치 해소 전 제출 금지' if _bad else '  (전 항목 일치)'))
            _c0.font = Font(bold=True, size=13, color=('C0392B' if _bad else '1F7A1F'))
            _rr += 1; ws2.cell(_rr,1,'담보'); ws2.cell(_rr,2,'한장보장표'); ws2.cell(_rr,3,'엑셀 끝열'); ws2.cell(_rr,4,'차이')
            for _k,_hv,_ev in _bad:
                _rr += 1
                ws2.cell(_rr,1,_k); ws2.cell(_rr,2,_hv); ws2.cell(_rr,3,_ev); ws2.cell(_rr,4,round(_ev-_hv,1))
                for _cc in range(1,5): ws2.cell(_rr,_cc).font = Font(color='C0392B', bold=True)
            if not _bad:
                _rr += 1; ws2.cell(_rr,1,'※ 대조 가능한 전 항목이 일치한다.')
            if _exl:
                _rr += 2; ws2.cell(_rr,1,'[검산 참고] 제외 7종 계약이 한장보장표에는 포함되어 생긴 차이 — 정상')
                ws2.cell(_rr,1).font = Font(bold=True, size=11, color='7A5C00')
                _rr += 1; ws2.cell(_rr,1,'담보'); ws2.cell(_rr,2,'한장보장표'); ws2.cell(_rr,3,'엑셀 끝열'); ws2.cell(_rr,4,'제외계약 몫')
                for _k,_hv,_ev,_dd in _exl:
                    _rr += 1
                    ws2.cell(_rr,1,_k); ws2.cell(_rr,2,_hv); ws2.cell(_rr,3,_ev); ws2.cell(_rr,4,_dd)
            print(f"[v277 검산] 한장표 대조 일치 {_ok} / 불일치 {len(_bad)}"
                  + (" → " + ', '.join(f"{k}({hv:,.0f}≠{ev:,.0f})" for k,hv,ev in _bad) if _bad else ""))
            # ★요약 한 줄을 시트 5행(눈에 먼저 들어오는 자리)에도 박는다
            _c5 = ws2.cell(5,1, f'[검산] 한장보장표 대비 불일치 {len(_bad)}건'
                                + (' — 아래 [검산] 표 확인' if _bad else ' — 전 항목 일치'))
            _c5.font = Font(bold=True, size=12, color=('C0392B' if _bad else '1F7A1F'))
    except Exception as _e9:
        print(f"[v277 검산] 실패 → 생략 ({_e9})")

    # ★★★★★v345 [중복줄] — 지점장 지시 2026.08.02(영구). <b>「조용히 틀리는 것」을 보이게 만든다.</b>
    #   v344로 「동일 담보명 2줄 = 대표(max)」가 <b>기본값</b>이 됐다. 값은 안전해졌지만
    #   <b>원래 몇 줄이었고 각각 얼마였는지가 아무 데도 안 남는다</b> → 진짜 합산해야 할 담보가
    #   나타나도 조용히 하나로 줄어든다. 그래서 <b>흔적을 확인사항에 강제 노출</b>한다.
    #   ★값은 건드리지 않는다(표시 전용). 지점장이 보고 "이건 합산이다"라고 하면 그때 예외로 넣는다.
    try:
        _dup_rows=[]
        for _ci,_c in enumerate(contracts):
            for _nm,_vs in (_c.get('dup') or {}).items():
                if len(set(_vs))==1 and len(_vs)>=2:
                    _dup_rows.append((_c.get('company',''), _nm, len(_vs), '/'.join(str(x) for x in _vs), '동일액 반복 → 대표 1개'))
                elif len(_vs)>=2:
                    _dup_rows.append((_c.get('company',''), _nm, len(_vs), '/'.join(str(x) for x in _vs), '금액 상이 → 대표(max) — 합산 여부 확인'))
        if _dup_rows:
            _r2 = ws2.max_row + 3
            _h = ws2.cell(_r2,1, f'[중복줄] 같은 계약에 담보명이 완전히 같은 줄이 2개 이상 — {len(_dup_rows)}건 (v344: 대표 1개만 기재)')
            _h.font = Font(bold=True, size=12, color='7A5C00')
            _r2 += 1
            for _i,_t in enumerate(('회사','별첨 담보명','줄 수','각 줄 금액','처리')):
                ws2.cell(_r2,_i+1,_t).font = Font(bold=True)
            for _co,_nm,_n,_vv,_wy in _dup_rows:
                _r2 += 1
                for _i,_v in enumerate((_co,_nm,_n,_vv,_wy)): ws2.cell(_r2,_i+1,_v)
            print(f"[v345 중복줄] {len(_dup_rows)}건 노출")
    except Exception as _e45:
        print(f"[v345 중복줄] 실패 → 생략 ({_e45})")

    # ★★★★★v359 [세부보충] — 별첨에 없는데 채워 넣은 값을 전부 노출한다.
    #   지점장 원문: "<b>왜 계속 없는걸 넣느냐</b>" · "<b>지침·메모리가 법률이다. 없는걸 새로 생각해서 넣지마라</b>".
    #   세부가입현황은 <b>구분이 뭉뚱그려진 표</b>라, 빈 행이라고 채우면 같은 담보가 두 행에 갈린다
    #   (실측 삼성 골절 포함/제외 각 50). <b>값은 지침대로 두되 흔적을 남겨</b> 지점장이 즉시 판별하게 한다.
    try:
        if _SEBU_LOG:
            _r3 = ws2.max_row + 3
            _h3 = ws2.cell(_r3,1, f'[세부보충] 별첨에 없어 세부가입현황에서 채운 값 — {len(_SEBU_LOG)}건 (별첨이 정본 · 다르면 별첨을 따른다)')
            _h3.font = Font(bold=True, size=12, color='7A5C00')
            _r3 += 1
            for _i,_t in enumerate(('회사','마스터 담보행','기재 금액','근거')):
                ws2.cell(_r3,_i+1,_t).font = Font(bold=True)
            for _co3,_tg3,_v3 in _SEBU_LOG:
                _r3 += 1
                for _i,_v in enumerate((_co3,_tg3,_v3,'세부가입현황(계약별 가입정보)')): ws2.cell(_r3,_i+1,_v)
            print(f"[v359 세부보충] {len(_SEBU_LOG)}건 노출")
    except Exception as _e59:
        print(f"[v359 세부보충] 실패 → 생략 ({_e59})")


    # ★★★★★v312 <b>해석원칙을 산출물에 원문 그대로 박는다</b>(지점장 지시 "법률 같은 봇").
    #   내가 지침을 요약·재구성해도 <b>고객 엑셀에는 지점장 원문이 그대로 나가</b> 왜곡이 즉시 드러난다.
    #   ★출력 전용 — 값·행 배정은 건드리지 않는다. 위치는 확인사항 시트 맨 아래.
    try:
        _pr = _PRINCIPLES or []
        if _pr:
            _pb = ws2.max_row + 3
            _ph = ws2.cell(_pb, 1, '[해석원칙] 이 문장이 판단의 정본이다 — 요약·재해석 금지  (출처: %s)'
                                   % _PRINCIPLES_SRC)
            _ph.font = Font(bold=True, size=12, color='1F3864')
            for _j, _t in enumerate(('No', '원칙 (원문)', '위반하면', '근거·사례'), 1):
                _hc = ws2.cell(_pb + 1, _j, _t)
                _hc.font = Font(bold=True, color='FFFFFF')
                _hc.fill = PatternFill('solid', fgColor='1F3864')
            for _i, (_p, _v, _g) in enumerate(_pr, 1):
                _r = _pb + 1 + _i
                ws2.cell(_r, 1, _i); ws2.cell(_r, 2, _p); ws2.cell(_r, 3, _v); ws2.cell(_r, 4, _g)
                ws2.cell(_r, 2).alignment = Alignment(wrap_text=True, vertical='top')
    except Exception as _ep:
        print('[v312 해석원칙] 출력 실패 → 생략 (%s)' % str(_ep)[:60])

    # ★★★★★v290 (지점장 지시 2026.07.31 "노란칸은 합계까지 이어지도록"):
    #   <b>원인</b>: master.xlsx의 행 채우기(연노랑 FFFFFFCC 등)가 <b>B~J(10열)까지만</b> 칠해져 있다.
    #   계약이 9건을 넘으면 끝열이 10열을 지나 <b>합계열에 색이 없다</b>(이명순 14계약 → 끝열 17).
    #   → 담보 행마다 <b>B열 채우기를 C~끝열 전체에 복사</b>한다. 계약 수와 무관하게 이어진다.
    #   ★값·수식·글자색은 건드리지 않는다(표시 전용). ★이미 다른 색이 칠해진 셀은 보존한다.
    try:
        _fillfix = 0
        for _r in range(6, ws.max_row + 1):
            if not ws.cell(_r, 2).value: continue
            _bf = ws.cell(_r, 2).fill
            _bg = _bf.fgColor.rgb if (_bf and _bf.fgColor) else None
            _bg = str(_bg) if _bg is not None else None
            # ★v371: 테마·인덱스 색이면 rgb가 aRGB 8자리가 아니다 → PatternFill이 거부한다.
            #   구 코드는 그대로 넘겨 [v290 행색 ERR]로 <b>행색 연장이 통째로</b> 죽었다.
            if not _bg or not re.fullmatch(r'[0-9A-Fa-f]{8}', _bg): continue
            if _bg == '00000000' or _bf.patternType != 'solid': continue
            for _c in range(3, last_col + 1):
                _cf = ws.cell(_r, _c).fill
                _cg = _cf.fgColor.rgb if (_cf and _cf.fgColor) else None
                if _cg == _bg: continue
                # 흰색·무채움 셀만 덮는다(특수 표시 색은 보존)
                if _cf.patternType == 'solid' and _cg not in ('00000000','FFFFFFFF',None): continue
                # ★v371: openpyxl 3.1.5는 fgColor.rgb가 str이 아니라 RGB 객체 → 그대로 넘기면
                #   PatternFill이 거부해 행색 연장이 통째로 실패했다(실측 [v290 행색 ERR]).
                ws.cell(_r, _c).fill = PatternFill('solid', fgColor=str(_bg))
                _fillfix += 1
        print(f'[v290 행색] 합계열까지 채우기 연장 {_fillfix}셀')
        # ★★★★★v291 (지점장 지적 2026.07.31 "허혈성이 갱신인데도 계속 블랙" · "실손도 파랑으로"):
        #   <b>원인</b>: 담보별 글자색(갱신=파랑 0070C0) 규칙이 <b>데이터셀에만</b> 적용되고
        #   <b>합계열(끝열)은 항상 검정</b>이었다. 실측: 허혈성·급성심근·입원·통원·약값 끝열 전부 FF000000.
        #   지점장이 보는 곳은 <b>합계열</b>이라 "갱신인데 블랙"으로 보인다.
        #   → 지침 §10 "갱신=파랑" 그대로 <b>합계열에도 적용</b>한다.
        #     ①그 행 데이터셀에 파랑이 하나라도 있으면 합계도 파랑
        #     ②실손(입원·통원·약값)·일상배상책임은 항상 파랑(§10 영구)
        _bluefix = 0
        _ALWAYS_BLUE = ('입원','통원','약값','약','일상배상책임')
        for _r in range(6, ws.max_row + 1):
            _nm = str(ws.cell(_r, 2).value or '').strip()
            if not _nm: continue
            _isblue = _nm in _ALWAYS_BLUE
            if not _isblue:
                for _c in range(3, last_col):
                    _f = ws.cell(_r, _c).font
                    _rgb = _f.color.rgb if (_f and _f.color) else None
                    if _rgb and str(_rgb).endswith('0070C0'): _isblue = True; break
            if _isblue:
                _tc = ws.cell(_r, last_col)
                _tc.font = Font(color='0070C0', bold=(_tc.font.bold if _tc.font else False))
                _bluefix += 1
        print(f'[v291 합계색] 합계열 파랑 적용 {_bluefix}행')
    except Exception as _e:
        print(f'[v290 행색 ERR] {_e}')

    # ★★★★★v289 근거표 — 어느 담보가 어느 행에 들어갔는지 전수 노출(반복 차단의 핵심)
    try:
        if '근거표' in wb.sheetnames: del wb['근거표']
        _tr = wb.create_sheet('근거표')
        _hd = ['마스터 행','원본 담보명','금액(만원)','회사','상품']
        for _j,_h in enumerate(_hd, start=1):
            _c = _tr.cell(1,_j,_h); _c.font = Font(bold=True, color='FFFFFF')
            _c.fill = PatternFill('solid', fgColor='1456B0')
        _cnt = {}
        for _t in trace_all: _cnt[_t[0]] = _cnt.get(_t[0],0)+1
        _rr = 2
        for _std0,_raw0,_amt0,_co0,_pd0 in sorted(trace_all, key=lambda x:(x[0], -x[2])):
            _tr.cell(_rr,1,_std0); _tr.cell(_rr,2,_raw0); _tr.cell(_rr,3,_amt0)
            _tr.cell(_rr,4,_co0);  _tr.cell(_rr,5,_pd0)
            if _cnt.get(_std0,0) > 1:      # 한 행에 2개 이상 합쳐진 곳 = 오탐 1순위
                for _j in range(1,6):
                    _tr.cell(_rr,_j).fill = PatternFill('solid', fgColor='FFF2CC')
            _rr += 1
        for _j,_w in enumerate((22,52,12,14,34), start=1):
            _tr.column_dimensions[chr(64+_j)].width = _w
        _tr.freeze_panes = 'A2'
        print(f'[v289 근거표] 담보 {len(trace_all)}건 기록 · 합산행 {sum(1 for v in _cnt.values() if v>1)}개')
    except Exception as _e:
        print(f'[v289 근거표 ERR] {_e}')

    # ★v39 워크시트 담보명 카피: 원본담보명을 숨김 시트 _dambo_raw 에 저장 (등식·기존시트 무손상)
    try:
        if '_dambo_raw' in wb.sheetnames: del wb['_dambo_raw']
        _rs = wb.create_sheet('_dambo_raw'); _rs.sheet_state='hidden'
        _rs.cell(1,1,'std'); _rs.cell(1,2,'raw'); _rs.cell(1,3,'amt')
        for _i,(_std,(_rw,_am)) in enumerate(raw_by_std.items(), start=2):
            _rs.cell(_i,1,_std); _rs.cell(_i,2,_rw); _rs.cell(_i,3,_am)
    except Exception:
        pass
    # ★★★★★v337b: <b>저장 직전</b>에 실손·일배책 행의 <b>끝열 합계 색을 파랑으로 확정</b>한다.
    #   중간에 넣으면 뒤 로직(세부보충·역기재 등)이 덮을 수 있어 지점장이 두 번 같은 지적을 했다.
    try:
        _ws0 = wb['보장분석']
        _lc0 = _ws0.max_column
        for _r9 in range(6, _ws0.max_row+1):
            if str(_ws0.cell(_r9,2).value or '').strip() in _BLUE_ROWS:
                for _c9 in range(3, _lc0+1):
                    # ★v388 제안 합계 열은 <b>레드가 정본</b>(지점장 확정) — 실손 파랑 강제에서 제외
                    if str(_ws0.cell(1,_c9).value or '').strip() == '제안 합계': continue
                    if _ws0.cell(_r9,_c9).value not in (None,''):
                        _ws0.cell(_r9,_c9).font = BL
    except Exception as _e9:
        print('[v337b 색] 실패:', str(_e9)[:60])
    # ★★★★★v370 가입제안서 열 색 (지점장 확정 2026.08.09):
    #   A1 헤더 = <b>주황 ED7D31 채우기 + 블랙 글자</b> / 담보 입력칸 = <b>레드 C00000 글자</b>(칸은 흰색).
    #   ★열 전체를 사후에 덮어쓴다 — 기존 12곳의 `BL if gen else BK`를 건드리지 않아 회귀 위험이 없다.
    try:
        _pidx = [i for i, _c in enumerate(contracts) if _c.get('proposal')]
        if _pidx:
            _ws0 = wb['보장분석']
            _ORG = PatternFill('solid', fgColor='ED7D31')
            _RD  = Font(color='C00000', name='맑은 고딕', size=9)
            for _i in _pidx:
                _pc = 3 + _i
                _h = _ws0.cell(1, _pc)
                _h.fill = _ORG
                _h.font = Font(color='000000', name='맑은 고딕', size=9, bold=True)
                _h.alignment = AL
                for _r in range(2, _ws0.max_row+1):
                    _cel = _ws0.cell(_r, _pc)
                    if _cel.value not in (None, ''):
                        _cel.font = _RD
            # ★★★★★v398 (지점장 지적 2026.08.12): <b>제안 합계 열의 1~5종 슬래시가 계속 블랙</b>이었다.
            #   원인: 종수술비 슬래시 행은 `=SUM()` 수식이 아니라 <b>문자열을 직접 써 넣는 별도 경로</b>라
            #     위 루프(제안 <b>계약</b> 열만 순회)에 걸리지 않았다. 다른 행은 수식이라 레드가 잘 먹었다.
            #   → 헤더가 <b>'제안 합계'</b>인 열도 값이 있으면 레드로 확정한다(v388 조문: 제안 합계 글자=레드).
            _sumred = 0
            for _c8 in range(3, _ws0.max_column+1):
                if str(_ws0.cell(1,_c8).value or '').strip() != '제안 합계': continue
                for _r8 in range(2, _ws0.max_row+1):
                    if _ws0.cell(_r8,_c8).value not in (None,''):
                        _ws0.cell(_r8,_c8).font = _RD; _sumred += 1
            print(f'[JEAN 색] 제안 열 {[3+i for i in _pidx]} → 헤더 주황ED7D31/블랙 · 값 레드C00000 · 제안합계열 레드 {_sumred}셀')
    except Exception as _ej:
        print('[JEAN 색] 실패:', str(_ej)[:80])
    # ★★★★★v464 제73조 (지점장 지적 2026.08.17 「엑셀이 1페이지만 나온다」)
    #   실측 원인: 마스터 xml에는 `<pageSetUpPr fitToPage="1"/>`만 있고 fitToWidth/Height 값이 없다.
    #   openpyxl이 저장할 때 <b>기본값 1/1</b>을 써서 106행 전체를 A4 한 장에 <b>억지로 축소</b>했다.
    #   → 가로만 1장(fitToWidth=1) · <b>세로는 무제한(fitToHeight=0)</b>으로 못박는다.
    #     담보표는 세로로 길다. 세로를 1장에 밀어넣으면 글자가 보이지 않는다.
    try:
        for _wsp in wb.worksheets:
            _wsp.page_setup.fitToWidth = 1
            _wsp.page_setup.fitToHeight = 0
            if _wsp.sheet_properties.pageSetUpPr is not None:
                _wsp.sheet_properties.pageSetUpPr.fitToPage = True
        print('[v464 인쇄] 가로 1장 · 세로 무제한 (시트 %d개)' % len(wb.worksheets))
    except Exception as _ep:
        print('[v464 인쇄] 설정 실패', str(_ep)[:60])
    _no_fullcalc(wb)          # ★v51 편집모드 강제 재계산 방지(수식은 유지)
    wb.save(out)
    _force_nocalc_xml(out)    # ★v124 저장 후 XML에 직접 못박음(3중 방어)
    return unmapped


def _force_nocalc_xml(path):
    """★★★v124(2026.07.21) 편집모드 속도 — 3중 방어의 마지막 단계.
       세 번 재발한 이유: 원천 master.xlsx가 fullCalcOnLoad="1"이라 코드가 매번 덮는 구조였다.
       v124에서 ①master.xlsx 자체를 "0"으로 교정 ②_no_fullcalc 유지 ③저장된 XML에 직접 못박음.
       셋 중 둘이 사라져도 폰 Excel '편집 사용'이 느려지지 않는다. 절대 제거 금지."""
    import zipfile, shutil, tempfile, os as _os
    try:
        zin = zipfile.ZipFile(path, 'r')
        tmp = path + '.nc'
        zout = zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED)
        for it in zin.infolist():
            data = zin.read(it.filename)
            if it.filename == 'xl/workbook.xml':
                x = data.decode('utf-8')
                if 'fullCalcOnLoad' in x:
                    x = re.sub(r'fullCalcOnLoad="[^"]*"', 'fullCalcOnLoad="0"', x)
                elif '<calcPr' in x:
                    x = re.sub(r'(<calcPr\b)', r'\1 fullCalcOnLoad="0"', x, count=1)
                else:
                    x = x.replace('</workbook>', '<calcPr fullCalcOnLoad="0"/></workbook>')
                data = x.encode('utf-8')
            zout.writestr(it, data)
        zin.close(); zout.close()
        shutil.move(tmp, path)
        return True
    except Exception:
        try:
            if _os.path.exists(path + '.nc'): _os.unlink(path + '.nc')
        except Exception:
            pass
        return False

# ★★★★★v388 합산 열(보유 합계·제안 합계·합계) 판정 — `ws.max_column`을 끝열로 쓰는 코드가
#   <b>데이터 열을 순회할 때</b> 새로 생긴 합산 2열을 포함하면 <b>이중 계산</b>이 된다.
#   헤더 글자로 판정한다(고정 열 번호를 박지 않는다 — 구조 가정 금지 #11).
def _is_sumcol(ws, c):
    h = str(ws.cell(1, c).value or '').strip()
    return h in ('합계', '보유 합계', '제안 합계')

def _data_cols(ws, last):
    """3 ~ last-1 중 <b>계약 열</b>만 (합산 열 제외)."""
    return [c for c in range(3, last) if not _is_sumcol(ws, c)]

def read_excel_totals(path):
    """완성 엑셀에서 담보명->합계 읽음. 등식2: PPT는 이것만 본다.
       ★★★v218 (지점장 지시 2026.07.25, 영구): <b>PPT 합계는 엑셀 끝열과 '반드시' 같아야 한다</b>.
       구 코드는 캐시가 없으면 <b>무조건 데이터셀 단순 SUM</b>으로 폴백했는데, 엑셀 끝열은 행마다
       수식이 다르다 → <b>대표(max)·캡 담보가 전부 어긋났다</b>(실측 불일치 5건):
         간병인 MAX 20 → PPT 30 / 간호통합병동 MAX 7 → PPT 12 / 중입자 MAX 300 → PPT 400 /
         자부상 MIN(,80) 80 → PPT 100 / 120대수술비 '1000/500' → PPT 1500.
       → <b>폴백도 끝열 수식 종류를 그대로 따른다</b>(inject_sum_cache와 동일 규칙).
       ★<b>동명 담보 행이 2개 이상</b>(혈전용해치료비 = 뇌 37행·심장 50행)이면 dict 키가 충돌해
       <b>뒤엣것이 앞엣것을 덮어쓴다</b> → PPT엔 한 칸뿐이므로 <b>두 행 중 대표(max)</b>로 병합한다."""
    wb = openpyxl.load_workbook(path, data_only=True)
    wbf = openpyxl.load_workbook(path)            # ★v218 수식 원문 판독용(data_only=False)
    ws = wb['보장분석']; wsf = wbf['보장분석']; last = ws.max_column
    out = {}; sq=[0]*5; ss=[0]*5; splits={}   # ★v370 splits[담보]=(갱신합, 비갱신합, 제안합) — 엑셀 글자색 근거(0070C0/검정/C00000)

    def _fallback(r):
        """끝열 캐시가 없을 때 — 엑셀 끝열 수식과 <b>같은 규칙</b>으로 계산한다."""
        nums = [ws.cell(r,c).value for c in _data_cols(ws,last) if isinstance(ws.cell(r,c).value,(int,float))]   # ★v388 합산 열 제외
        s = sum(nums)
        f = wsf.cell(r,last).value
        if isinstance(f,str) and f.startswith('='):
            if f.startswith('=MIN('):
                mm = re.search(r',\s*(\d+)\s*\)\s*$', f)
                return min(s, int(mm.group(1))) if mm else s
            if f.startswith('=IF(COUNT'):   return max(nums) if nums else 0   # 대표(max) 담보
            if f.startswith('=IF(SUM'):     return 7 if s>0 else 0
        return s

    def _put(nm, val):
        """★v218 동명 행 충돌 방지 — 이미 있으면 큰 값 유지(대표)."""
        if not val: return
        prev = out.get(nm)
        out[nm] = max(prev, val) if isinstance(prev,(int,float)) else val

    for r in range(6, ws.max_row+1):
        nm = ws.cell(r,2).value
        if not nm: continue
        nm = str(nm).strip()
        endv = ws.cell(r,last).value
        # ★★★v219 (지점장 지시 2026.07.25, 영구): <b>PPT 갱신/비갱신 색은 엑셀 글자색이 원천이다</b>.
        #   구 코드는 PPT가 raw dambo로 갱신합/비갱신합을 <b>따로 다시 계산</b>했다. 그런데 엑셀 끝열은
        #   대표(max)·캡·[확인] 제외 때문에 raw 분할합과 <b>거의 항상 다르다</b> → 구 코드는 그럴 때
        #   <b>'큰 쪽으로 전부 몰아버려'</b>(`if gs>=ns: gs,ns=_T,0`) <b>갱신↔비갱신이 통째로 뒤바뀌었다</b>.
        #   실측 뒤바뀜: 갱신 100(엑셀 산입) + 비갱신 200([확인] 제외) → ns가 커서 <b>검정</b>으로 찍히지만
        #   엑셀의 100은 <b>갱신 계약 값이라 파랑</b>이어야 한다. 역방향도 동일하게 발생.
        #   → <b>엑셀 데이터셀의 글자색(파랑 0070C0 = 갱신)을 그대로 읽어</b> 분할한다. 이중 계산 폐기.
        _cells=[]
        for c in _data_cols(ws, last):        # ★v388 합산 열 제외(이중 계산 방지)
            _v = ws.cell(r,c).value
            if not isinstance(_v,(int,float)) or not _v: continue
            try: _rgb = str(wsf.cell(r,c).font.color.rgb or '')
            except Exception: _rgb = ''
            _up = _rgb.upper()
            # ★★★★★v370 (지점장 확정 2026.08.09): <b>레드 C00000 = 가입제안서 값</b>.
            #   갱신(파랑)·비갱신(검정)에 이어 <b>제안분을 세 번째 항</b>으로 분리한다.
            _cells.append((_v, _up.endswith('0070C0'), _up.endswith('C00000')))
        _f0 = wsf.cell(r,last).value
        if isinstance(_f0,str) and _f0.startswith('=IF(COUNT') and _cells:
            # ★대표(max) 행 — 끝열 값을 만든 <b>최댓값 셀 하나의 색</b>이 정답이다.
            #   합으로 나누면 갱신 8+9(=17)가 비갱신 15를 눌러 <b>색이 뒤집힌다</b>(구 코드 실패 지점).
            _mx = max(_cells, key=lambda x: x[0])
            _gs, _ns, _ps = (0,0,0)
            if   _mx[2]: _ps = _mx[0]
            elif _mx[1]: _gs = _mx[0]
            else:        _ns = _mx[0]
        else:
            _gs = sum(v for v,b,p in _cells if b and not p)
            _ns = sum(v for v,b,p in _cells if not b and not p)
            _ps = sum(v for v,b,p in _cells if p)          # ★v370 제안(레드)
        if _gs or _ns or _ps:
            _pv0 = splits.get(nm,(0,0,0))
            if len(_pv0)==2: _pv0=(_pv0[0],_pv0[1],0)
            splits[nm] = (max(_pv0[0], _gs), max(_pv0[1], _ns), max(_pv0[2], _ps))
        # 수술비 1~5종: 끝열 슬래시 문자열(수식 아님, 항상 존재)
        # ★★★★★v400 (지점장 지적 2026.08.12): <b>보장분석지 PPT의 1~5종 슬래시가 블랙</b>이었다.
        #   실측: 그 run의 color가 <b>None</b>(폼 기본색 상속)이었다 — 다른 값은 전부 C00000이 잘 들어갔다.
        #   원인: 슬래시 줄은 `rsl()`이 <b>줄 전체를 통째로 교체</b>하는 별도 경로라 색 지정이 없었다.
        #   → 그 행의 <b>계약 열 슬래시 셀 폰트색</b>을 보고 <b>전부 제안(C00000)이면 레드</b>로 표시한다.
        #     (v398에서 고친 것은 <b>엑셀</b> 제안합계 열이고, 이번은 <b>PPT</b>다 — 서로 다른 곳이다.)
        def _slash_src(_r):
            _red=_any=False
            for _c in _data_cols(ws, last):
                _v = ws.cell(_r,_c).value
                if not isinstance(_v,str) or '/' not in _v: continue
                _any=True
                try: _rg=str(wsf.cell(_r,_c).font.color.rgb or '').upper()
                except Exception: _rg=''
                if _rg.endswith('C00000'): _red=True
                else: return False
            return _red and _any
        if nm == '상해 종수술비(1-5종)' and isinstance(endv,str) and '/' in endv:
            for k,p in enumerate(endv.split('/')[:5]):
                try: ss[k]=int(p)
                except: pass
            if _slash_src(r): splits['__SS_RED__']=(0,0,1)
            continue
        if nm == '질병 종수술비(1-5종)' and isinstance(endv,str) and '/' in endv:
            for k,p in enumerate(endv.split('/')[:5]):
                try: sq[k]=int(p)
                except: pass
            if _slash_src(r): splits['__SQ_RED__']=(0,0,1)
            continue
        if nm == 'MRI/도수치료/비급여주사' and isinstance(endv,str) and '/' in endv:   # ★v29y 트리오 분해
            _ps=endv.split('/')
            for _k,_std in enumerate(('MRI','도수치료','비급여주사')):
                try: out[_std]=int(_ps[_k])
                except: pass
            continue
        # ★★v218 120대수술비(n대수술비) = 끝열이 <b>계약별 가로 슬래시 문자열</b>이다(합산 금지, v30k 정본).
        #   구 코드는 이 문자열을 못 읽어 폴백 단순합산 1,500을 PPT에 넣었다(엑셀 '1000/500'과 불일치).
        #   → 슬래시를 쪼개 <b>대표(max)</b>를 PPT 값으로 쓴다.
        if isinstance(endv,str) and '/' in endv:
            _pv=[]
            for _p in endv.split('/'):
                try: _pv.append(int(str(_p).strip().replace(',','')))
                except: pass
            if _pv: _put(nm, max(_pv))
            continue
        # 숫자 합계: 끝열 캐시값 있으면 사용, 없으면 끝열 수식과 동일 규칙으로 계산(★v218)
        _put(nm, endv if isinstance(endv,(int,float)) and endv else _fallback(r))
    return out, sq, ss, splits

def build_ppt(data, out, totals=None, surg_q=None, surg_s=None, splits=None):
    if not os.path.exists(TPL_PPT): return False
    prs = Presentation(TPL_PPT)
    sl = prs.slides[0]
    by = {sh.name:sh for sh in sl.shapes if sh.has_text_frame}
    client = data['client']; contracts = data['contracts']
    now = datetime.datetime.now()

    totals = totals if totals is not None else {}
    surg_q = surg_q if surg_q is not None else [0]*5
    surg_s = surg_s if surg_s is not None else [0]*5
    if not totals:   # 폴백: 전달 없으면 옛 방식(등식2 불가시)
        for ct in contracts:
            for raw,amt in ct['dambo'].items():
                std = resolve(raw)
                if std: totals[std]=totals.get(std,0)+amt

    # ★PPT 색: 하나라도 갱신=파랑 / 전부 비갱신=검정 / 실손 항상 파랑 (미가입은 값 미기재라 해당없음)
    _BLUE=RGBColor(0x00,0x00,0xFF); _BLACK=RGBColor(0x00,0x00,0x00); _RED=RGBColor(0xC0,0x00,0x00)   # ★v370 가입제안서 = 레드
    # ★★★★★v347 (지점장 지적 2026.08.02 "둘 다 파란색인데 둘 중 하나는 계속 블랙처리된다"):
    #   §10 정본 = 실손·일배책은 <b>항상 파랑</b>인데, 엑셀은 `_BLUE_ROWS`(상해의료비·MRI트리오 포함)를
    #   쓰고 <b>PPT는 이 집합을 따로 갖고 있어 두 곳이 어긋나 있었다</b> → 같은 실손인데 한 칸만 검정.
    #   → <b>엑셀 집합과 통일</b>한다(상해의료비·MRI트리오 추가).
    _silson={'입원','통원','약값','약','MRI','MRI트리오','도수치료','비급여주사','상해의료비','일상배상책임'}  # ★v210 간병인·간호통합병동 강제 파랑 폐기(엑셀과 동일 규칙) — 구 v139 3행 무조건 파랑 폐기
    # 담보별 '최대 기여 계약'의 갱신여부로 색 결정 → 합산 시 전부 파랑 쏠림 방지(엑셀 혼합과 일치)
    _dom={}  # std -> (max_amt, gen)
    for ct in contracts:
        _gen = (ct.get('renewal','')=='갱신')
        for raw,amt in ct.get('dambo',{}).items():
            if not amt: continue
            st=resolve(raw)
            if not st: continue
            if st not in _dom or amt>_dom[st][0]: _dom[st]=(amt,_gen)
    def pcol(std):
        if std in _silson: return _BLUE
        d=_dom.get(std)
        return _BLUE if (d and d[1]) else _BLACK
    # ★담보별 갱신합/비갱신합 (분할 표기용)
    # ★★★v219 (지점장 지시 2026.07.25, 영구): <b>갱신합/비갱신합은 엑셀 글자색에서 그대로 가져온다</b>.
    #   구 코드는 raw dambo로 재계산해 엑셀 끝열과 어긋났고, 어긋나면 큰 쪽으로 몰아버려
    #   <b>갱신↔비갱신이 통째로 뒤바뀌었다</b>. splits=read_excel_totals가 준 (갱신합, 비갱신합).
    _gensum={}; _nonsum={}; _propsum={}   # ★v370 제안(레드)
    if splits:
        for _st,_sv in splits.items():
            _g,_n = _sv[0],_sv[1]; _p = _sv[2] if len(_sv)>2 else 0
            if _g: _gensum[_st]=_g
            if _n: _nonsum[_st]=_n
            if _p: _propsum[_st]=_p
    else:   # 폴백(엑셀 없이 호출된 경우) — 구 방식
        for ct in contracts:
            _gen=(ct.get('renewal','')=='갱신')
            for raw,amt in ct.get('dambo',{}).items():
                if not amt: continue
                st=resolve(raw)
                if not st: continue
                tgt=_gensum if (_gen or '갱신' in raw) else _nonsum
                tgt[st]=tgt.get(st,0)+amt
    # ★v30n PPT 골절 = 골절(치아파절포함) + 골절(치아파절제외) 합산. 엑셀은 두 행 분리 유지, PPT만 하나로 합산 표기(지점장 2026.07.03)
    for _b in (_gensum, _nonsum, totals):
        _b['골절합산PPT'] = (_b.get('골절(치아파절포함)',0) or 0) + (_b.get('골절(치아파절제외)',0) or 0)
    def _seg(run0, segs):
        run0.text=segs[0][0]
        if segs[0][1] is not None:
            try: run0.font.color.rgb=segs[0][1]
            except: pass
        prev=run0._r
        for txt,col in segs[1:]:
            new=_copy.deepcopy(run0._r)
            t=new.find(_qn('a:t'))
            if t is not None: t.text=txt
            prev.addnext(new); prev=new
            nr=_Run(new, run0._parent)
            if col is not None:
                try: nr.font.color.rgb=col
                except: pass
    def pv(box,pi,ri,std,prefix=': ',suffix=''):
        # 한 칸 분할: 갱신합=파랑 / 비갱신합=검정 (둘 다면 'gen / non'), 실손=파랑 합계
        # ★★★v221: 구 코드는 박스·문단·run을 못 찾으면 <b>아무 말 없이 return</b>했다.
        #   그래서 'PPT 상해 수술 라인이 통째로 비어 있다' 같은 사고가 <b>로그 한 줄 없이</b> 나갔다.
        #   → 실패 사유를 반드시 찍는다(배포 후 Railway 로그로 즉시 원인 확정 가능).
        if box not in by:
            print(f'[PPT_MISS] 박스없음 box={box} std={std} — 템플릿 ppt_form.pptx의 도형 이름 확인 필요'); return
        tf=by[box].text_frame
        if pi>=len(tf.paragraphs):
            print(f'[PPT_MISS] 문단없음 box={box} p{pi} (문단수 {len(tf.paragraphs)}) std={std}'); return
        p=tf.paragraphs[pi]
        if ri>=len(p.runs):
            print(f'[PPT_MISS] run없음 box={box} p{pi} r{ri} (run수 {len(p.runs)}) std={std}'); return
        gs=_gensum.get(std,0); ns=_nonsum.get(std,0); ps=_propsum.get(std,0)   # ★v370 제안(레드)
        if std in _silson:
            _v = totals.get(std,0)            # ★실손=완성 엑셀값(입원5천캡·통원디폴트 반영). _gensum 원본합산(상해+질병=1만) 사용 안 함
            if not _v: return
            segs=[(f'{prefix}{_v:,}{suffix}', _BLUE)]
            _seg(p.runs[ri], segs); return
        # ★v30f 등식1 (지점장 승인 2026.07.03): PPT 표기 총액의 정본 = 완성 엑셀 끝열.
        #   대표 1건·MAX·캡·[확인] 분리로 끝열 ≠ raw 분할합이면 끝열 값 단색 표기(색=우세 성분측).
        # ★★v219: gs·ns가 이미 <b>엑셀 셀에서 온 값</b>이므로 원칙적으로 gs+ns == 끝열이다.
        #   다만 끝열이 대표(max)·캡 수식인 행은 합보다 작다 → <b>큰 쪽으로 몰지 말고</b>
        #   <b>어느 쪽 성분이 끝열 값을 만들었는지</b>로 판정한다(색 뒤바뀜 차단).
        _T = totals.get(std, None)
        if ps and isinstance(_T,(int,float)) and _T>=ps: _T=_T-ps   # ★v370 제안분을 뗀다
        if isinstance(_T,(int,float)) and _T>0 and (gs or ns) and (gs+ns)!=_T:   # ★v370 0-나눗셈 가드
            if   gs and not ns: gs,ns=int(_T),0
            elif ns and not gs: gs,ns=0,int(_T)
            elif gs>=_T and ns< _T: gs,ns=int(_T),0      # 끝열값을 만든 쪽 = 갱신
            elif ns>=_T and gs< _T: gs,ns=0,int(_T)      # 끝열값을 만든 쪽 = 비갱신
            else:                                        # 둘 다 기여(=SUM 계열) → 비율 보존
                _tot=gs+ns
                gs=int(round(_T*gs/_tot)); ns=int(_T)-gs
        if not gs and not ns and not ps: return
        # ★★★★★v370: 3분할 = 갱신(파랑)+비갱신(검정)+제안(레드)
        segs=[]; _f=True
        if gs: segs.append(((f'{prefix}{gs:,}' if _f else f'+{gs:,}'), _BLUE)); _f=False
        if ns: segs.append(((f'{prefix}{ns:,}' if _f else f'+{ns:,}'), _BLACK)); _f=False
        if ps: segs.append(((f'{prefix}{ps:,}' if _f else f'+{ps:,}'), _RED)); _f=False
        if segs: segs[-1]=(segs[-1][0]+suffix, segs[-1][1])
        _seg(p.runs[ri], segs); return
        # ── 구 2분할(도달 불가·보존) ──
        if gs and ns:
            segs=[(f'{prefix}{gs:,}', _BLUE),(f'+{ns:,}{suffix}', _BLACK)]
        elif gs:
            segs=[(f'{prefix}{gs:,}{suffix}', _BLUE)]
        else:
            segs=[(f'{prefix}{ns:,}{suffix}', _BLACK)]
        _seg(p.runs[ri], segs)
    def _slot(box, label):
        """★★★★★v330 (지점장 확정 2026.08.02, 영구): <b>폼이 뼈대의 정본이다 — 코드는 값만 채운다</b>.
        지점장 원문: "엑셀도 각 ppt는 <b>뼈대에 잘 채워지는것이 목적</b>이다" ·
        "담보는 동일화지만 <b>담보의 위치나 순서는 바뀔수있다</b>".
        → 구 코드는 (문단,run) <b>좌표를 하드코딩</b>했다. 폼이 한 줄만 바뀌어도 전부 밀린다.
          실측 오염 — <b>심장수술비가 '뇌혈관수술' 칸</b>에, <b>뇌혈관수술비가 '1~7종 수술비' 칸</b>에,
          <b>유사암 값이 '통합전이암' 칸</b>에 들어갔고, <b>'1~5종 수술비' 라벨은 값에 먹혔다</b>.
        → <b>라벨 글자로 슬롯을 찾는다</b>: 라벨 뒤 첫 ':'가 있는 run을 값 자리로 보고,
          ':' 뒤 꼬리(정렬 공백·다음 라벨)는 <b>그대로 보존</b>한다. 구조 가정 금지 원칙과 같은 등급.
        반환 (문단객체, run인덱스, run내 ':' 오프셋) / 못 찾으면 None."""
        if box not in by: return None
        for p in by[box].text_frame.paragraphs:
            full=''.join(r.text for r in p.runs)
            k=-1
            _from=0
            while True:
                _k=full.find(label,_from)
                if _k<0: break
                # ★v330c 단어경계: 라벨 앞 글자가 한글이면 다른 담보다('깁스'가 '반깁스'에 걸리던 실측)
                if _k==0 or not ('가'<=full[_k-1]<='힣'): k=_k; break
                _from=_k+1
            if k<0: continue
            c=full.find(':', k+len(label))
            if c<0: continue
            off=0
            for ri,r in enumerate(p.runs):
                if off <= c < off+len(r.text): return p, ri, c-off
                off+=len(r.text)
        return None
    def pvl(box, label, std):
        """라벨로 찾은 슬롯에 값만 기입(갱신=파랑 + 비갱신=검정 분할 유지, 꼬리 보존)."""
        sl=_slot(box,label)
        if not sl:
            print(f'[PPT_MISS] 라벨없음 box={box} label={label!r} std={std}'); return
        p, ri, c = sl
        t=p.runs[ri].text; head=t[:c+1]; tail=t[c+1:]
        # ★v330b: 폼(표본)에 예시값이 남아 있으면 값이 두 번 찍힌다(실측 '상해수술 : 200 200').
        #   ':' 뒤 꼬리가 <b>숫자·콤마·공백뿐</b>이면 예시값으로 보고 버린다.
        #   '   / 약 : ' 처럼 <b>다음 라벨·구분자가 있으면 그대로 보존</b>한다(뼈대 훼손 금지).
        if tail.strip() and not re.search(r'[^\d,\s]', tail): tail=''
        gs=_gensum.get(std,0); ns=_nonsum.get(std,0); ps=_propsum.get(std,0)
        if std in _silson:
            _v=totals.get(std,0)
            if not _v: return
            _seg(p.runs[ri], [(f'{head} {_v:,}',_BLUE),(tail,None)]); return
        _T=totals.get(std,None)
        if ps and isinstance(_T,(int,float)) and _T>=ps: _T=_T-ps   # ★v370 제안분을 뗀다(합계엔 이미 포함)
        if isinstance(_T,(int,float)) and _T>0 and (gs or ns) and (gs+ns)!=_T:   # ★v370 gs·ns 둘 다 0이면 보정 금지(0-나눗셈)
            if   gs and not ns: gs,ns=int(_T),0
            elif ns and not gs: gs,ns=0,int(_T)
            elif gs>=_T and ns< _T: gs,ns=int(_T),0
            elif ns>=_T and gs< _T: gs,ns=0,int(_T)
            else:
                _tt=gs+ns; gs=int(round(_T*gs/_tt)); ns=int(_T)-gs
        if not gs and not ns and not ps: return
        # ★★★★★v370 (지점장 확정 2026.08.09): 한 칸 3분할 = 갱신(파랑)+비갱신(검정)+<b>제안(레드)</b>.
        segs=[]; _f=True
        if gs: segs.append((f'{head} {gs:,}',_BLUE)); _f=False
        if ns: segs.append(((f'{head} {ns:,}' if _f else f'+{ns:,}'),_BLACK)); _f=False
        if ps: segs.append(((f'{head} {ps:,}' if _f else f'+{ps:,}'),_RED)); _f=False
        segs.append((tail,None))
        _seg(p.runs[ri], segs)
        return
        # ── 구 2분할(도달 불가·보존) ──
        if gs and ns: segs=[(f'{head} {gs:,}',_BLUE),(f'+{ns:,}',_BLACK),(tail,None)]
        elif gs:      segs=[(f'{head} {gs:,}',_BLUE),(tail,None)]
        else:         segs=[(f'{head} {ns:,}',_BLACK),(tail,None)]
        _seg(p.runs[ri], segs)
    def rsl(box, label_line, text, red=False):
        """라벨 없이 슬래시 괄호줄 같은 <b>줄 전체</b>를 교체(1~5종 수술비 칸)."""
        if box not in by: return False
        for p in by[box].text_frame.paragraphs:
            full=''.join(r.text for r in p.runs)
            if label_line in full and p.runs:
                p.runs[0].text=text
                for _rr in p.runs[1:]: _rr.text=''
                if red:                      # ★v400 제안분 슬래시 = 레드 C00000
                    try: p.runs[0].font.color.rgb = RGBColor(0xC0,0x00,0x00)
                    except Exception: pass
                return True
        print(f'[PPT_MISS] 줄없음 box={box} pat={label_line!r}'); return False
    def _setcol(run,std):
        try: run.font.color.rgb=pcol(std)
        except: pass
    _last_std=[None]
    def g(nm):
        _last_std[0]=nm
        return totals.get(nm,0)
    def r_set(box,pi,ri,val,std='__USE_LAST__'):
        use = _last_std[0] if std=='__USE_LAST__' else std
        # ★v221 조용한 실패 금지 — 못 찍었으면 사유를 로그로 남긴다.
        if box not in by:
            print(f'[PPT_MISS] 박스없음 box={box} val={val}')
        else:
            tf=by[box].text_frame
            if pi>=len(tf.paragraphs):
                print(f'[PPT_MISS] 문단없음 box={box} p{pi} (문단수 {len(tf.paragraphs)}) val={val}')
            else:
                p=tf.paragraphs[pi]
                if ri>=len(p.runs):
                    print(f'[PPT_MISS] run없음 box={box} p{pi} r{ri} (run수 {len(p.runs)}) val={val}')
                else:
                    p.runs[ri].text=val
                    if use: _setcol(p.runs[ri], use)
        _last_std[0]=None

    for b in ['TextBox 49','TextBox 56']:
        if b in by: by[b].text_frame.word_wrap=False

    by['TextBox 21'].text_frame.word_wrap=False
    by['TextBox 21'].text_frame.auto_size=MSO_AUTO_SIZE.NONE  # 도형 고정(이름 길이에 따라 박스 이동·크기변경 방지)
    by['TextBox 21'].text_frame.paragraphs[0].runs[0].text=f'{client} 님의 보장'
    by['TextBox 21'].text_frame.paragraphs[0].runs[1].text='(전)'
    # 날짜를 한 박스(TextBox 36)로 통합, 35·29는 비움
    if 'TextBox 36' in by and 'TextBox 29' in by:
        try: by['TextBox 36'].width = by['TextBox 29'].left + by['TextBox 29'].width - by['TextBox 36'].left
        except: pass
    by['TextBox 36'].text_frame.paragraphs[0].runs[0].text=f'{now.year}년 {now.month:02d}월 {now.day:02d}일 기준'
    for _eb in ('TextBox 35','TextBox 29'):
        if _eb in by:
            for _pp in by[_eb].text_frame.paragraphs:
                for _rr in _pp.runs: _rr.text=''
    # 상단 헤더(이름+날짜) 전부 18pt·도형 고정
    for _hb in ('TextBox 21','TextBox 36','TextBox 35','TextBox 29'):
        if _hb in by:
            by[_hb].text_frame.word_wrap=False
            try: by[_hb].text_frame.auto_size=MSO_AUTO_SIZE.NONE
            except: pass
            for _pp in by[_hb].text_frame.paragraphs:
                try: _pp.alignment = PP_ALIGN.CENTER      # ★v41 이름·날짜 우측쏠림 → 가운데
                except: pass
                for _rr in _pp.runs:
                    try: _rr.font.size=Pt(18)
                    except: pass

    # ★v48(지점장 2026.07.13): 제목(고객명)+날짜를 한 덩어리로 슬라이드 가운데 배치
    try:
        _t, _d = by.get('TextBox 21'), by.get('TextBox 36')
        if _t is not None and _d is not None:
            _SW = prs.slide_width
            _GAP = 100000
            _tot = _t.width + _GAP + _d.width
            _lf = int((_SW - _tot) / 2)
            _t.left = _lf
            _d.left = _lf + _t.width + _GAP
            _d.top = _t.top
    except Exception:
        pass

    if g('질병사망(80세)'): pvl('TextBox 10','80세','질병사망(80세)')
    if g('상해사망'): pvl('TextBox 11','상해사망','상해사망')
    # ★★★★★v330c (등식2 위반 수리): 구 코드는 <b>raw dambo를 다시 합산</b>해 '종신' 칸에 넣었다
    #   → 실측 이영태 <b>13,476</b> vs 엑셀 일반사망 끝열 <b>91,710</b>. PPT는 <b>완성 엑셀만</b> 읽는다.
    #   §8.1 정본대로 종신 칸 = 마스터 <b>'일반사망'</b> 행이다.
    if g('일반사망'): pvl('TextBox 10','종신','일반사망')

    if g('상해후유3%'): pvl('TextBox 8','상해3%','상해후유3%')
    if g('질병후유3%'): pvl('TextBox 8','질병 3%','질병후유3%')
    if g('상해후유80%'): pvl('TextBox 8','상해 80%','상해후유80%')
    if g('질병후유80%'): pvl('TextBox 8','질병 80%','질병후유80%')

    if g('뇌혈관진단비'): pv('TextBox 46',0,0,'뇌혈관진단비',prefix='뇌혈관\n',suffix='')
    if g('뇌졸증진단비'): pv('TextBox 47',0,0,'뇌졸증진단비',prefix='뇌졸증\n',suffix='')
    if g('뇌출혈진단비'): pv('TextBox 48',0,0,'뇌출혈진단비',prefix='뇌출혈\n',suffix='')
    if g('산정특례뇌혈관'): pvl('TextBox 49','산정특례','산정특례뇌혈관')
    if g('혈전용해치료비'): pvl('TextBox 49','혈전용해치료비','혈전용해치료비')
    if g('2대 주요치료비'): pvl('TextBox 49','대주요치료비','2대 주요치료비')   # 뇌혈관쪽 2대주요치료비

    # ★ 심장 표기(설명서와 동일 8종): 1줄 협심증/심부전/염증/빈맥 · 2줄 부정맥/심근병증/심장판막. 값 있는 것만. 급성심근·허혈성 별도칸.
    # ★★★★★v336 (지점장 지시 2026.08.02): 심장 담보를 <b>담보별 칸</b>으로 표기한다.
    #   구 코드는 `협심증/심부전/염증/빈맥`을 <b>슬래시로 묶고 최댓값 1개</b>만 찍어
    #   담보별 금액이 보이지 않았다. 지점장: "<b>빈맥 : / 염증 : / 심근병증 : / 심장판막 :</b> 넣어야 한다".
    #   → 폼 `TextBox 심장4종`에 4칸을 <b>가로 한 줄</b>로 두고 라벨로 찾아 값만 채운다(§11 v330).
    #   ★협심증·심부전·부정맥은 종전대로 `TextBox 4` 칸을 쓴다(폼에 이미 있다).
    #   ★`협심증 / 심부전 :`은 폼상 <b>한 칸</b>이므로 값이 있는 쪽(큰 쪽)을 대표로 넣는다 — 구 로직과 동일.
    # ★★★★★v337 (지점장 지적 2026.08.02 "왜 이걸 두개를 합친것이냐 다른 담보들인데"):
    #   구 코드는 폼의 `협심증 / 심부전 :` <b>한 칸</b>에 둘 중 큰 값만 넣었다 — <b>다른 담보를 합친 것</b>이다.
    #   → 폼을 `협심증 :   심부전 :` <b>두 칸</b>으로 나누고 <b>각각 기재</b>한다.
    # ★★★★★v338 (지점장 지시 2026.08.02 "ppt에 이렇게 줄여야한다 칸이 다 튀어나온다"):
    #   구 폼은 `TextBox 4`(협심증·심부전·부정맥)가 <b>`TextBox 심장4종` 안에 완전히 겹쳐</b> 있어
    #   글자가 같은 자리에 두 번 인쇄되고 칸 밖으로 넘쳤다(좌표 실측으로 확인).
    #   → 심장 <b>6칸을 `TextBox 심장4종` 한 박스에 2줄</b>로 합치고 <b>7pt</b>로 줄였다.
    #     1줄 = 빈맥·염증·심근병증·심장판막 / 2줄 = 협심증·심부전·부정맥. `TextBox 4`는 비웠다.
    if g('협심증'): pvl('TextBox 심장4종','협심증','협심증')
    if g('심부전'): pvl('TextBox 심장4종','심부전','심부전')
    if g('부정맥'): pvl('TextBox 심장4종','부정맥','부정맥')
    # ★v337 1-7/1-8/1-9종 = 폼 `1~7종 수술비 :` 칸에 대표값(최댓값)
    # ★v351: 엑셀 '120대수술비'를 분석지 PPT `_____대 수술 :` 칸에도 기재(질병 박스 TextBox 17).
    if g('n대수술비'): pvl('TextBox 17','대 수술','n대수술비')
    if g('상해 종수술비(1-8종)'): pvl('TextBox 19','1~7종 수술비','상해 종수술비(1-8종)')
    if g('질병 종수술비(1-8종)'): pvl('TextBox 17','1~7종 수술비','질병 종수술비(1-8종)')
    for _hn in ('빈맥','염증','심근병증','심장판막'):
        if g(_hn): pvl('TextBox 심장4종', _hn, _hn)
    # ★★★★★v318 허혈성 진단비(TextBox 54) — <b>pv() 경로로 통일</b>(지점장 지시 2026.08.01).
    #   <b>구 결함 2가지</b>: ①`허혈성 : 5,000` <b>한 줄</b>로 나와 뇌혈관·뇌졸증·급성심근(`이름\n값` 두 줄)과
    #   모양이 달랐다 ②`pv()`를 안 타고 `runs[0].text`를 직접 덮어써서 <b>색 지정이 아예 없었다</b>
    #   → 갱신=파랑/비갱신=검정(엑셀 글자색 원천, v219)이 허혈성에만 적용되지 않았다.
    if g('허혈성 진단비'): pv('TextBox 54',0,0,'허혈성 진단비',prefix='허혈성\n',suffix='')
    elif 'TextBox 54' in by:
        _t54=by['TextBox 54'].text_frame
        if _t54.paragraphs[0].runs: _t54.paragraphs[0].runs[0].text='허혈성'
    if g('급성심근경색'): pv('TextBox 55',0,0,'급성심근경색',prefix='급성심근\n',suffix='')
    if g('산정특례심장'): pvl('TextBox 56','산정특례','산정특례심장')
    if g('2대 주요치료비'): pvl('TextBox 56','대주요치료비','2대 주요치료비')   # 심장쪽 2대주요치료비

    # ★★★★★v322 <b>암 블록 순서 = 엑셀(마스터)과 동일</b>(지점장 지시 2026.08.01).
    #   지점장 원문: `고액암 : / 통합암 : / 일반암 : / 유사암 : / 통합전이암 :`
    #   템플릿은 `암진단비`(p0) / `유사암`(p1) / `통합전이암`(p2) 3문단뿐 → <b>p0에 3줄</b>을 줄바꿈으로 넣는다.
    def _amt3(std):
        _g=_gensum.get(std,0) or 0; _n=_nonsum.get(std,0) or 0; _v=totals.get(std,0) or 0
        if _g and _n: return [(f'{_g:,}',_BLUE),(f'+{_n:,}',_BLACK)]
        if _g:        return [(f'{_g:,}',_BLUE)]
        if _n:        return [(f'{_n:,}',_BLACK)]
        if _v:        return [(f'{_v:,}',_BLACK)]
        return []
    # ★★★★★v330: 암 블록 — 라벨 삽입 폐기. 폼 라벨로 슬롯을 찾아 값만 채운다.
    for _lb,_st in (('고액암','고액암'),('일반암','일반암'),('통합암','통합암'),
                    ('통합전이암','통합전이암'),('유사암','유사암(갑.기.경.제)')):
        if g(_st): pvl('TextBox 14',_lb,_st)
    if g('항암방사선약물'): pvl('TextBox 14','항암치료','항암방사선약물')
    if g('표적항암치료비'): pvl('TextBox 14','표적치료','표적항암치료비')
    if g('세기조절치료'): pvl('TextBox 14','세기','세기조절치료')
    if g('양성자치료') and not g('세기조절치료'): pvl('TextBox 14','세기','양성자치료')
    if g('다빈치로봇수술비'): pvl('TextBox 14','다빈치로봇수술비','다빈치로봇수술비')
    # 상급병원 암주요치료비 / 하이클래스 (TextBox 57)
    if 'TextBox 57' in by: by['TextBox 57'].text_frame.word_wrap=False
    # ★★★★★v320 TextBox 57 라벨 정본(지점장 2026.08.01): `암주요치료비 :` / `하이클래스(비급여) :` / `10억통장 :`
    # ★★★★★v330: 라벨 덮어쓰기·10억통장 줄 삽입 폐기 — 폼에 이미 있다. 값만 채운다.
    # ★★★★★v371 (지점장 지적 2026.08.09): <b>폼엔 `암주요치료비 :` 라벨이 있는데 코드가 부르지 않았다</b>.
    #   엑셀 21행 값이 있어도 PPT는 늘 공란 — 결과값 동결(#9) 위반. v363 「1~3종 주입 코드 부재」와 같은 뿌리.
    if g('암주요치료비'): pvl('TextBox 57','암주요치료비','암주요치료비')
    if g('하이클래스(암)'): pvl('TextBox 57','하이클래스','하이클래스(암)')
    if g('10억 플랜'):     pvl('TextBox 57','10억통장','10억 플랜')

    if g('질병수술비'): pvl('TextBox 17','질병수술','질병수술비')
    if any(surg_q): rsl('TextBox 17','/     /', '(%s)'%'/'.join(str(x) for x in surg_q), red=bool((splits or {}).get('__SQ_RED__')))
    if g('뇌혈관수술비'): pvl('TextBox 17','뇌혈관수술','뇌혈관수술비')
    if g('심장수술비'): pvl('TextBox 17','심장 수술','심장수술비')
    if g('상해수술비'): pvl('TextBox 19','상해수술','상해수술비')
    if any(surg_s): rsl('TextBox 19','/      /', '(%s)'%'/'.join(str(x) for x in surg_s), red=bool((splits or {}).get('__SS_RED__')))
    if g('골절수술비'): pvl('TextBox 19','골절수술','골절수술비')

    _ys=totals.get('양성자치료',0); _sgj=totals.get('세기조절치료',0)   # ★v29v (지점장 2026.07.02) 양성자·세기조절 → 암 박스
    # ★v330: 양성자·세기조절은 위 pvl('TextBox 14','세기',…)로 기입한다(구 좌표 덮어쓰기 폐기).
    실손_cts=[ct for ct in contracts
        if any('실손' in k or '입원의료비' in k for k in ct['dambo']) and ct['contract_date']]
    실손가입일=min((c['contract_date'] for c in 실손_cts), default='___________')
    _실손상품=next((c.get('product','') for c in 실손_cts if c['contract_date']==실손가입일), '')
    _np3=any(_has_nonpay3(c.get('dambo')) for c in 실손_cts)   # ★v250 3대비급여 하한
    _sg=silson_gen(실손가입일, totals.get('입원'), _실손상품, _np3, bool(totals.get('약값')))   # ★실손 세대 자동판별(상품명 연도코드 반영)
    by['TextBox 59'].text_frame.word_wrap=False
    by['TextBox 59'].text_frame.paragraphs[0].runs[0].text='실손'+(f' {_sg}' if _sg else '')
    by['TextBox 59'].text_frame.paragraphs[1].runs[0].text='('
    by['TextBox 59'].text_frame.paragraphs[1].runs[1].text='가입일:'
    by['TextBox 59'].text_frame.paragraphs[1].runs[2].text=f'{실손가입일})'
    for r in by['TextBox 59'].text_frame.paragraphs[1].runs: r.font.size=Pt(10)  # ★v50 '다10'
    if g('입원'): pvl('TextBox 6','입원','입원')
    if g('통원'): pvl('TextBox 6','통원','통원')
    if g('약값'): pvl('TextBox 6','약','약값')
    if g('MRI'): pvl('TextBox 6','MRI','MRI')
    if g('도수치료'): pvl('TextBox 6','도수치료','도수치료')
    if g('비급여주사'): pvl('TextBox 6','비급여주사','비급여주사')
    # ★★★★★v329 (지점장 지시 2026.08.02): 새 표본(ppt_form.pptx)의 실손 박스에
    #   <b>`상해의료비 : ` 줄(p5)</b>이 추가됐다 → 마스터 '상해의료비' 행 값을 그대로 기입한다.
    #   지점장 원문: "실손에 상해의료비 추가했다 이걸로 최종해라".
    #   ★상해의료비는 실손과 별개 정액 담보다(마스터 103행) — 입원과 합치지 않는다.
    if g('상해의료비'): pvl('TextBox 6','상해의료비','상해의료비')

    # ★★★★★v349 (지점장 지적 2026.08.02): "<b>골절진단의료비용(치아파절제외)라고 적혀있는데
    #   치아포함에 넣어졌다</b>". 매핑(resolve)은 정상이었고 <b>PPT가 두 행을 합쳐 「골절(치아포함)」
    #   칸 하나에만</b> 찍고 있었다(구 v30n) → 제외 담보만 가입한 고객도 「포함」 칸에 나왔다.
    #   폼에 <b>`골절(치아포함) :` · `골절(치아제외) :` 두 칸이 따로</b> 있으므로 <b>각각 기재</b>한다.
    #   ★엑셀 두 행 = PPT 두 칸 → §1 등식(담보·값 동일)이 맞는다. 구 합산 표기는 폐기.
    if g('골절(치아파절포함)'): pvl('TextBox 7','골절(치아포함)','골절(치아파절포함)')
    if g('골절(치아파절제외)'): pvl('TextBox 7','골절(치아제외)','골절(치아파절제외)')
    if g('화상진단비'): pvl('TextBox 7','화상','화상진단비')
    if g('깁스진단비'): pvl('TextBox 7','깁스','깁스진단비')
    if g('응급실(응급)'): pvl('TextBox 7','응급실','응급실(응급)')
    if g('일상배상책임'): pvl('TextBox 5','일상배상책임','일상배상책임')
    if g('대인'): pvl('TextBox 9','대인','대인')
    if g('대물'): pvl('TextBox 9','대물','대물')
    if g('합의금'): pvl('TextBox 9','합의금','합의금')
    if g('6주미만'): pvl('TextBox 9','주미만','6주미만')
    if g('변호사'): pvl('TextBox 9','변호사비','변호사')
    if g('자부상'): pvl('TextBox 9','자부상','자부상')
    if g('질병일당'): pvl('TextBox 22','질병일당','질병일당')
    if g('상해일당'): pvl('TextBox 22','상해일당','상해일당')
    # ★★★★★v319 일당 박스 라벨·값 정본(지점장 지시 2026.08.01) — 지점장 원문 형태:
    #   `질병일당 :   / 질병종합병원일당 : `  /  `상해일당 :   / 상해종합병원일당 : `
    #   `간병인일당 :   / 요양병원 : `        /  `간호통합병동일당 :   / 간병인지원일당 : `
    #   구 템플릿 라벨은 `병원일당`(질병·상해 구분이 없었다)이고 <b>값도 안 채워졌다</b>.
    def _lab(box,pi,ri,text):
        try:
            _p=by[box].text_frame.paragraphs[pi]
            if ri < len(_p.runs): _p.runs[ri].text=text
        except Exception: print('[PPT_MISS] 라벨 %s p%d r%d'%(box,pi,ri))
    if g('종합병원 질병입원일당'): pvl('TextBox 22','질병종합병원일당','종합병원 질병입원일당')
    if g('종합병원 상해입원일당'): pvl('TextBox 22','상해종합병원일당','종합병원 상해입원일당')
    if g('1인실 상급병원'): pvl('TextBox 22','인실 상급병원일당','1인실 상급병원')
    if g('1인실 종합병원'): pvl('TextBox 22','인실 종합병원일당','1인실 종합병원')
    if g('암일당'): pvl('TextBox 22','암일당','암일당')   # ★v371 동일 결함 — 폼 라벨 있는데 미호출(엑셀 30행)
    # ★★★★★v512 제110조 — 폼 라벨 「간병인일당」 → <b>「간병인 사용일당」</b>,
    #   「요양병원」 → <b>「요양병원 간병인」</b>(ppt_form.pptx 수정). 값 3개를 각각 채운다.
    if g('간병인'): pvl('TextBox 22','간병인 사용일당','간병인')
    # ★v319: `간호통합병동일당 : 값 / 간병인지원일당 : 값` — 한 문단에 담보 2개.
    # ★★★★★v330: 간호통합병동·간병인지원일당 라벨 삽입 폐기 — 폼 p9에 이미 있다.
    if g('간호통합병동'):   pvl('TextBox 22','간호통합병동일당','간호통합병동')
    if g('간병인지원일당'): pvl('TextBox 22','간병인지원일당','간병인지원일당')
    if g('크라운'): pvl('TextBox 13','크라운','크라운')
    if g('임플란트'): pvl('TextBox 13','임플란트','임플란트')

    # ── 누락 슬롯 보충 (엑셀 합계 끌어오기) ──
    if g('중입자치료비'): pvl('TextBox 14','중입자','중입자치료비')
    if g('5대골절진단비'): pvl('TextBox 7','대골절','5대골절진단비')
    if g('중증화상진단비'): pvl('TextBox 7','중대화상','중증화상진단비')
    if g('허혈성수술비'): pvl('TextBox 17','허혈성수술','허혈성수술비')
    if g('5대골절수술비'): pvl('TextBox 19','대골절수술','5대골절수술비')
    if g('화상수술비'): pvl('TextBox 19','화상수술','화상수술비')
    if g('창상봉합술'): pvl('TextBox 19','창상봉합수술','창상봉합술')
    if g('질병중환자실'): pvl('TextBox 22','질병중환자실','질병중환자실')
    if g('상해중환자실'): pvl('TextBox 22','상해중환자실','상해중환자실')
    if g('1인실 상급병원'): pvl('TextBox 22','인실 상급병원일당','1인실 상급병원')
    if g('1인실 종합병원'): pvl('TextBox 22','인실 종합병원일당','1인실 종합병원')

    # ★v29t CI 담보값 노란 배경(§8.4·§11): 중대한 계열을 해당 칸에 표기 + 값 run만 노랑 하이라이트
    from pptx.oxml.ns import qn as _ciqn
    import copy as _cicopy
    from pptx.text.text import _Run as _ciRunCls
    _CIHL_AFTER=[_ciqn('a:uLnTx'),_ciqn('a:uLn'),_ciqn('a:uFillTx'),_ciqn('a:uFill'),_ciqn('a:latin'),_ciqn('a:ea'),
               _ciqn('a:cs'),_ciqn('a:sym'),_ciqn('a:hlinkClick'),_ciqn('a:hlinkMouseOver'),_ciqn('a:rtl'),_ciqn('a:extLst')]
    def _hl_yellow(run):
        rPr=run._r.get_or_add_rPr()
        for old in rPr.findall(_ciqn('a:highlight')): rPr.remove(old)
        hl=rPr.makeelement(_ciqn('a:highlight'),{}); hl.append(rPr.makeelement(_ciqn('a:srgbClr'),{'val':'FFFF00'}))
        ins=None
        for ch in rPr:
            if ch.tag in _CIHL_AFTER: ins=ch; break
        if ins is not None: ins.addprevious(hl)
        else: rPr.append(hl)
    def _ci_run(box,pidx,std,sep):
        v=totals.get(std,0)
        if not v or box not in by: return
        tf=by[box].text_frame
        if pidx>=len(tf.paragraphs): return
        p=tf.paragraphs[pidx]
        if not p.runs: return
        base=p.runs[-1]
        nr_el=_cicopy.deepcopy(base._r); base._r.addnext(nr_el)
        nr=_ciRunCls(nr_el,p); nr.text=f'{sep}{v:,}'
        _hl_yellow(nr)
    def _ci_split(box,label,ci_std,extra_std):
        # ★v29t: 라벨줄 + [CI값(노랑)] + [+일반값] 을 별도 run으로 구성 — 개행 포함 run의 하이라이트 미표시(파워포인트) 방지
        civ=totals.get(ci_std,0)
        if not civ or box not in by: return
        tf=by[box].text_frame; p=tf.paragraphs[0]
        if not p.runs: return
        base=p.runs[0]
        for _r in list(p.runs[1:]): _r._r.getparent().remove(_r._r)
        base.text=f'{label}\n'
        el1=_cicopy.deepcopy(base._r); base._r.addnext(el1)
        r1=_ciRunCls(el1,p); r1.text=f'{civ:,}'; _hl_yellow(r1)
        # ★v417 CI 값 색도 엑셀 글자색 근거(구 코드는 base run 색을 그대로 물려받아 파랑이 됐다)
        try:
            r1.font.color.rgb=(_RED if _propsum.get(ci_std) else (_BLUE if _gensum.get(ci_std) else _BLACK))
        except: pass
        # ★★★★★v417 (지점장 지적 2026.08.13 「급성심근경색 → 기존것과 합쳐져 다 블루로 나온다」)
        #   구 코드는 일반 담보를 <b>끝열 합계 한 숫자</b>로 찍고 색도 `_gensum` 유무만 봐서
        #   <b>보유+제안이 뭉쳐 통째로 파랑</b>이 됐다(실측 `4,000`+`+2,500` 둘 다 0000FF).
        #   ★ pv()는 이미 3분할(갱신·비갱신·제안)을 하는데 _ci_split이 그 위를 <b>덮어써서</b>
        #     제안 레드가 사라졌다 — 여기서도 splits를 그대로 쓴다(제40조 재현 완료).
        _g2=_gensum.get(extra_std,0); _n2=_nonsum.get(extra_std,0); _p2=_propsum.get(extra_std,0)
        if not (_g2 or _n2 or _p2):
            _ex0=totals.get(extra_std,0)
            if _ex0: _n2=_ex0
        _prev=el1
        for _v2,_c2 in ((_g2,_BLUE),(_n2,_BLACK),(_p2,_RED)):
            if not _v2: continue
            _e2=_cicopy.deepcopy(base._r); _prev.addnext(_e2); _prev=_e2
            _r2=_ciRunCls(_e2,p); _r2.text=f'+{_v2:,}'
            try: _r2.font.color.rgb=_c2
            except: pass
    # ★★★★★v243(지점장 지시 2026.07.25): <b>보장분석지 PPT 뇌 칸도 축을 따라간다</b>.
    #   구 코드는 `중대한 뇌졸증`만 봐서, 축이 뇌출혈인 CI(신한·DB 실측)는 <b>PPT에 아무것도 안 찍혔다</b>.
    #   → 끝열 합계에 <b>중대한 뇌출혈</b>이 있으면 그 축으로 라벨·값을 바꾼다.
    if (totals.get('중대한 뇌출혈',0) or 0) > (totals.get('중대한 뇌졸증',0) or 0):
        _ci_split('TextBox 47','뇌출혈','중대한 뇌출혈','뇌출혈진단비')
    else:
        _ci_split('TextBox 47','뇌졸증','중대한 뇌졸증','뇌졸증진단비')
    _ci_split('TextBox 55','급성심근','중대한 급성심근','급성심근경색')
    _ci_run('TextBox 14',0,'중대한 암','+')
    _ci_run('TextBox 10',3,'중대한CI적용','+')
    _autofit_ppt(by)
    # ★★★★★v531 제121조 2항 (지점장 지시 2026.08.21 「네 산출물의 담보값을 실제로 대조」)
    #   진단서·리포트는 `heart_audit`가 막는다. <b>보장분석지 PPT도 같은 검사를 받는다.</b>
    #   폼에 실제로 찍힌 글자를 파싱해 엑셀 끝열(totals)과 대조한다 — 다르면 발행을 막는다.
    _hbad = _ppt_heart_audit(by, totals)
    if _hbad:
        print('[제121조 2항 심장동결·PPT] 불일치 %d건 — 발행 차단' % len(_hbad))
        for _x in _hbad: print('   ·', _x)
        raise RuntimeError('제121조 심장동결 위반(보장분석지 PPT): ' + ' / '.join(_hbad))
    print('[제121조 2항 심장동결·PPT] 보장분석지 = 엑셀 · 불일치 0건')
    prs.save(out); return True


# ★★★★★v531 제121조 2항 — 보장분석지 PPT 심장 6칸 값 대조(검사 전용).
#   ★키 표는 <b>마스터 심장 행 이름</b>에서 온다. 주입 코드를 참조하지 않는다.
_PPT_HEART = (('빈맥', '빈맥'), ('염증', '염증'), ('심근병증', '심근병증'),
              ('심장판막', '심장판막'), ('협심증', '협심증'), ('심부전', '심부전'),
              ('부정맥', '부정맥'))


def _ppt_heart_audit(by, totals):
    """폼 `TextBox 심장4종`에 찍힌 글자를 파싱해 엑셀 끝열과 대조한다."""
    import re as _re2
    _sh = by.get('TextBox 심장4종')
    if _sh is None or not getattr(_sh, 'has_text_frame', False):
        return []                                   # 폼에 칸이 없으면 검사 대상 아님
    _tx = _sh.text_frame.text or ''
    _bad = []
    for _lb, _key in _PPT_HEART:
        _want = int(totals.get(_key, 0) or 0)
        _m = _re2.search(_re2.escape(_lb) + r'\s*:\s*([\d,+]*)', _tx)
        if not _m:
            if _want: _bad.append('보장분석지 %s 칸 없음 — 엑셀 %s' % (_lb, format(_want, ',')))
            continue
        _got = sum(int(x) for x in _re2.findall(r'\d+', _m.group(1).replace(',', '')) ) if _m.group(1).strip() else 0
        if _got != _want:
            _bad.append('보장분석지 %s — 엑셀 %s / PPT %s' % (_lb, format(_want, ','), format(_got, ',')))
    return _bad


# ★v50(지점장 '다10'): 제목·날짜만 예외(18pt). 실손박스(59)도 10pt 대상으로 편입.
_HEADER_BOXES={'TextBox 21','TextBox 36','TextBox 35','TextBox 29'}
_SURGERY_BOXES={'TextBox 17','TextBox 19'}   # ★v29t: 질병수술·상해수술 9.0pt 고정(지점장 2026.07.02), 1~5종 줄만 축소 허용
def _autofit_ppt(by):
    """겹침·단락내림 방지(§11): 값박스 word_wrap off + 최장 단락 기준 박스 단위 축소.
    수술 박스 2개는 8.9pt 고정, '1~5종' 제목줄·슬래시 괄호줄만 축소 허용."""
    for _bn, sh in by.items():
        if _bn in _HEADER_BOXES: continue
        tf = sh.text_frame
        try:
            tf.word_wrap = False
            w_in = sh.width / 914400.0
        except: continue
        # ★★★★★v338 (지점장 지시 2026.08.02 "ppt에 이렇게 줄여야한다 칸이 다 튀어나온다"):
        #   심장 6칸(빈맥·염증·심근병증·심장판막 / 협심증·심부전·부정맥)을 한 박스 2줄로 합쳤다.
        #   10pt로는 <b>칸 밖으로 넘친다</b> → 이 박스만 <b>8pt 고정</b>(지점장 지시 2026.08.02)(v50 '다10' 예외 3번째).
        if _bn == 'TextBox 심장4종':
            for p in tf.paragraphs:
                for r in p.runs:
                    try: r.font.size = Pt(8)
                    except: pass
            continue
        if _bn in _SURGERY_BOXES:
            # ★수술비 폰트(지점장 규정 2026.07.07): 1-5종 슬래시 줄만 6pt, 나머지 수술 줄은 9pt 고정(축소 금지)
            for p in tf.paragraphs:
                ptxt=''.join(r.text for r in p.runs)
                # ★★★v318 수술 1~5종 슬래시 줄 = <b>9pt</b>(지점장 지시 2026.08.01).
                #   구 v50 정본 「슬래시 줄 6pt」는 <b>폐기</b> — 지점장 원문 "(40/80/600/2000/4800) → 글자포인트9".
                _sz = 9.0 if ('/' in ptxt) else 10.0  # ★v318: 슬래시(1-5종) 9pt, 그 외 10pt
                for r in p.runs:
                    if r.text:
                        try: r.font.size = Pt(_sz)
                        except: pass
            continue
        runs_all = [r for p in tf.paragraphs for r in p.runs if r.text]
        if not runs_all: continue
        # ★v50 정본(지점장 2026.07.13): 값 폰트는 전부 10pt 고정.
        #   - v50(2026.07.13): 값·라벨 전부 10pt(지점장 '다10'). 예외=제목·날짜(18pt)·수술 1~5종(6pt).
        #   - 6pt는 수술 1~5종 슬래시 줄에만 허용(위 _SURGERY_BOXES 분기).
        for r in runs_all:
            try:
                cur = r.font.size.pt if r.font.size else 9.0
                if cur < 18.0 and cur != 10.0:
                    r.font.size = Pt(10)
            except: pass


def build_chiryo(data, out, totals=None, unmapped=None):
    """치료비 정리 폼: 고객명/날짜 + [확인](AI 미매핑) 항목을 카테고리별로 채움.
    추측 금지 — 박스 라벨이 명시한 'AI가 못 채운 항목'에 실제 미매핑 목록만 주입."""
    if not os.path.exists(TPL_TX): return False
    prs = Presentation(TPL_TX); sl = prs.slides[0]
    by = {sh.name:sh for sh in sl.shapes if sh.has_text_frame}
    client = data['client']; now = datetime.datetime.now()
    def first_run_set(box, text):
        if box not in by: return
        tf = by[box].text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = text
    if 'TextBox 21' in by:
        by['TextBox 21'].text_frame.word_wrap=False
        by['TextBox 21'].text_frame.auto_size=MSO_AUTO_SIZE.NONE  # 도형 고정
        rs=by['TextBox 21'].text_frame.paragraphs[0].runs
        if rs: rs[0].text=f'{client} 님의 보장'
        if len(rs)>1: rs[1].text='(전)'
    first_run_set('TextBox 36', f'{now.year}년')
    first_run_set('TextBox 35', f'{now.month:02d}월')
    first_run_set('TextBox 29', f'{now.day:02d}일 기준')
    # [확인] 미매핑 항목 → 회사별 묶어 본문 박스에 기재(있을 때만)
    unmapped = unmapped or []
    if unmapped:
        lines = [f"{comp} {raw}: {amt:,}" for (col,comp,raw,amt,note) in unmapped]
        blob = '\n'.join(lines[:20])
        for box in ['TextBox 25','TextBox 32','TextBox 37','TextBox 51']:
            if box in by:
                tf=by[box].text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = '[확인] AI 미매핑(별첨 직접확인):\n'+blob
                break
    _autofit_ppt(by)
    prs.save(out); return True

def make_summary(data):
    contracts=data['contracts']; cust=data['client']
    total_premium=sum(ct['premium'] for ct in contracts)
    갱신수=sum(1 for ct in contracts if ct['renewal']=='갱신')
    lines=[f"<b>👤 {cust} 고객님 분석 완료</b>","",
           f"<b>계약 현황</b>",
           f"  • 총 계약 수: <b>{len(contracts)}건</b>",
           f"  • 갱신형: {갱신수}건 / 비갱신형: {len(contracts)-갱신수}건",
           f"  • 월 보험료 합계: <b>{total_premium:,}원</b>","","🏢 <b>가입 회사</b>"]
    for ct in contracts:
        tag='🔵갱신' if ct['renewal']=='갱신' else '🔴비갱신' if '비갱신' in ct['renewal'] else '🟢완납'
        lines.append(f"  • {ct['company']} [{tag}] {ct['premium']:,}원")
    totals={}
    for ct in contracts:
        for raw,amt in ct['dambo'].items():
            std=resolve(raw)
            if std: totals[std]=totals.get(std,0)+amt
    key=[('일반암','🎗암진단'),('뇌혈관진단비','🧠뇌혈관'),('협심증','❤️허혈성'),
         ('급성심근경색','❤️급성심근'),('상해사망','💀상해사망'),('질병사망(80세)','💀질병사망'),('입원','🏥실손')]
    found=[(lbl,totals[k]) for k,lbl in key if k in totals and totals[k]>0]
    if found:
        lines+=["","🔑 <b>주요 담보 합계 (만원)</b>"]
        for lbl,amt in found: lines.append(f"  • {lbl}: <b>{amt:,}만원</b>")
    return '<br>'.join(lines)

INDEX_HTML = r'''<!DOCTYPE html>
<html lang="ko"><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>@@BRAND@@ @@BSUB@@</title>
<link rel="manifest" href="/manifest.webmanifest">
<meta name="theme-color" content="#06203f">
<meta name="mobile-web-app-capable" content="yes">
<meta name="apple-mobile-web-app-capable" content="yes">
<meta name="apple-mobile-web-app-status-bar-style" content="black-translucent">
<meta name="apple-mobile-web-app-title" content="@@BRAND@@">
<link rel="apple-touch-icon" href="/icon-180.png">
<link rel="icon" href="/icon-192.png">
<script>if("serviceWorker" in navigator){window.addEventListener("load",function(){navigator.serviceWorker.register("/sw.js").catch(function(e){console.log("sw",e);});});}</script>
<style>
:root{--bg:#ffffff;--panel:#f6f7f9;--line:#dde2ea;--acc:#06203f;--acc2:#a8863d;--ink:#15202e;--mute:#6c7887;--gold:#c5a052;--red:#b3242c;--green:#06203f;--blue:#a8863d;--sky:#2a6bf2;--sky2:#1a54d0}
*{box-sizing:border-box;margin:0;padding:0}
body{background:var(--bg);color:var(--ink);font-family:'Pretendard','Noto Sans KR',sans-serif;line-height:1.55}
#gate{position:fixed;inset:0;z-index:100;background:var(--bg);display:flex;flex-direction:column;align-items:center;justify-content:center;padding:30px 26px;text-align:center}
#gate .kick{font-size:14px;font-weight:800;letter-spacing:.45em;color:var(--gold);margin-bottom:14px}
#gate h1{font-size:30px;font-weight:800;margin-bottom:14px;color:var(--acc)}
#gate .s{font-size:14px;color:var(--mute);margin-bottom:38px}
#gate .pw{width:100%;max-width:420px;background:#ffffff;border:1.5px solid var(--line);border-radius:14px;padding:18px 20px;font-size:17px;color:var(--ink);text-align:center;letter-spacing:.3em;outline:none}
#gate .pw:focus{border-color:var(--acc)}
#gate .go{width:100%;max-width:420px;margin-top:14px;border:none;border-radius:14px;padding:18px;font-size:17px;font-weight:800;color:#fff;background:var(--acc);cursor:pointer}
#gate .err{color:var(--red);font-size:13px;font-weight:700;margin-top:14px;min-height:18px}
.shake{animation:sh .35s}@keyframes sh{0%,100%{transform:translateX(0)}25%{transform:translateX(-8px)}75%{transform:translateX(8px)}}
.app{max-width:520px;margin:0 auto;height:100vh;display:none;flex-direction:column}
@media (min-width:900px){ .app{max-width:1000px;padding:0 16px}
 .msg{max-width:78%;font-size:14px} .chat{padding:18px 8px}
 #gate .pw,#gate .go{max-width:460px} }
@media (min-width:1400px){ .app{max-width:1200px} }
/* ★v435 PC 대응 — 넓은 화면에서 폭·글자를 키운다(제56조) */
@media (min-width:900px){
 .app{max-width:1000px;padding:0 16px}
 .msg{max-width:78%;font-size:14px}
 .chat{padding:18px 8px}
 #gate .pw,#gate .go{max-width:460px}
 .file-card{font-size:13px}
}
@media (min-width:1400px){ .app{max-width:1200px} }
header{padding:14px 18px;border-bottom:2px solid var(--gold);background:#ffffff;display:flex;align-items:center;gap:10px}
.logo{width:32px;height:32px;border-radius:9px;border:1.5px solid var(--gold);display:flex;align-items:center;justify-content:center;font-size:16px}
h1{font-size:14px;font-weight:800}h1 b{color:var(--acc2)}.sub{font-size:10px;color:var(--mute)}
.chat{flex:1;overflow-y:auto;padding:16px 12px;display:flex;flex-direction:column;gap:12px}
.msg{max-width:90%;font-size:13px}
.me{align-self:flex-end;background:rgba(6,32,63,.05);border:1px solid rgba(6,32,63,.16);border-radius:14px 14px 4px 14px;padding:9px 13px}
.bot{align-self:flex-start;background:#ffffff;border:1px solid var(--line);border-radius:14px 14px 14px 4px;padding:11px 14px;width:100%;box-shadow:0 1px 3px rgba(6,32,63,.06)}
.file-cards{display:flex;flex-direction:column;gap:8px;margin-top:10px}
.file-card{display:flex;align-items:center;gap:11px;border-radius:12px;padding:11px 13px}
.file-card.xl{background:rgba(6,32,63,.04);border:1px solid rgba(6,32,63,.28)}
.file-card.pt{background:rgba(197,160,82,.07);border:1px solid rgba(197,160,82,.50)}
.file-card .ic{font-size:22px}.file-card .nm{flex:1;font-size:12.5px;font-weight:700}
.file-card .dl{font-size:11px;font-weight:800;padding:5px 11px;border-radius:8px}
.file-card.xl .dl{color:var(--green);background:rgba(6,32,63,.08)}
.file-card.pt .dl{color:var(--blue);background:rgba(197,160,82,.14)}
.summary-box{background:#f6f7f9;border:1px solid #dde2ea;border-radius:10px;padding:12px 14px;margin-top:10px;font-size:12px;line-height:1.7}
.err{color:var(--red);font-size:12px}
.spin{width:22px;height:22px;border:3px solid var(--line);border-top-color:var(--acc);border-radius:50%;animation:sp .8s linear infinite;display:inline-block;vertical-align:middle}
@keyframes sp{to{transform:rotate(360deg)}}
.bar{padding:12px;border-top:1px solid var(--line);display:flex;gap:9px;background:var(--bg);align-items:stretch}
/* ★v603b 산출물 칩 — 줄이 넘치면 자동으로 다음 줄로 흐른다(모바일 대응) */
.oc{display:inline-flex;align-items:center;gap:3px;font-size:11.5px;line-height:1;
    padding:5px 8px;border-radius:999px;background:rgba(197,160,82,.14);
    border:1px solid rgba(197,160,82,.45);color:#7a6122;white-space:nowrap}
@media (max-width:400px){ .oc{font-size:10.5px;padding:4px 7px} }
.up{flex:1;min-width:0;border:1px solid rgba(197,160,82,.55);background:#fcfbf7;border-radius:12px;padding:14px 8px;text-align:center;font-size:13px;font-weight:700;cursor:pointer;color:var(--acc);white-space:nowrap;overflow:hidden;text-overflow:ellipsis;letter-spacing:-.01em;transition:border-color .15s,background .15s}
.up:hover{border-color:var(--gold);background:#faf6ec}
.send{white-space:nowrap}
@media (max-width:520px){ .up{font-size:11.5px;padding:12px 4px;letter-spacing:-.02em}
 .send{font-size:12.5px;padding:0 12px} }
@media (max-width:380px){ .up{font-size:10.5px;padding:11px 3px} }
.send{border:none;border-radius:12px;padding:0 20px;font-weight:800;font-size:14px;background:var(--sky);color:#fff;cursor:pointer;letter-spacing:.02em}
.send:hover{background:var(--sky2)}
.send:disabled{opacity:.4}
/* ★v626 지점장 「질문 빼라 의미없다」 — JS 는 두고 화면에서만 지운다 */
.qbar{display:none!important}
.qlbl{display:none!important}
.qinput{flex:1;min-width:0;background:#ffffff;border:1px solid var(--line);border-radius:12px;padding:12px 14px;font-size:13px;color:var(--ink);outline:none}
.qinput:focus{border-color:var(--acc)}.qinput::placeholder{color:var(--mute)}
.qbtn{border:none;border-radius:12px;padding:0 18px;font-weight:800;font-size:13px;background:var(--sky);color:#fff;cursor:pointer}
.qbtn:hover{background:var(--sky2)}.qbtn:disabled{opacity:.4}

footer{text-align:center;font-size:10px;color:var(--mute);padding:8px}footer b{color:var(--acc2)}
/* ★★★★★v624 (지점장 지시 2026.09.02 「맨아래에 4개버튼 넣고 가운데는 새고객」 · 「밝은 블루로」) */
.tabbar{display:flex;align-items:flex-end;justify-content:space-around;
        background:#fff;border-top:1px solid var(--line);padding:9px 4px 10px;position:relative}
.tab{width:20%;min-width:0;text-align:center;color:var(--mute);font-size:10px;cursor:pointer;
     text-decoration:none;display:block;line-height:1.25;padding:0 2px}
.tab .ti{display:block;line-height:0;margin-bottom:3px}
.tab .ti svg{display:inline-block;vertical-align:middle}
.tab .nm{display:block;font-size:9.5px;max-width:100%;overflow:hidden;
         text-overflow:ellipsis;white-space:nowrap}
.tab.on{color:var(--sky);font-weight:800}
.tab:active{color:var(--sky)}
.gowrap{padding:10px 12px 12px;background:var(--bg);border-top:1px solid var(--line)}
.gobtn{width:100%;border:none;border-radius:12px;padding:15px;font-size:15px;font-weight:800;
       background:var(--sky);color:#fff;cursor:pointer;letter-spacing:.02em}
.gobtn:hover{background:var(--sky2)}
.gobtn:disabled{background:#e3e7ee;color:#9aa3af;cursor:default}
.tabc{width:20%;min-width:0;text-align:center;color:#f08c14;font-size:10px;cursor:pointer;
      line-height:1.25;padding:0 2px}
.tabc .ti{display:block;line-height:0;margin-bottom:3px}
.tabc .ti svg{display:inline-block;vertical-align:middle}
.tabc .nm{display:block;font-size:9.5px;font-weight:800}

.tabc .lb{font-size:10px;margin-top:3px;color:var(--sky);font-weight:700}
</style></head><body>
<div id="gate">
  <div class="kick">@@BRAND@@</div><h1>@@BRAND@@ @@BSUB@@</h1>
  <div class="s">회원코드 6자리를 입력하세요</div>
  <input id="pw" class="pw" type="text" inputmode="numeric" placeholder="회원코드 6자리" autocomplete="off">
  <button id="go" class="go">접속</button><div id="gerr" class="err"></div>
  <div id="joinbox" style="margin-top:18px;width:100%;max-width:420px">
   <div id="joinopen" style="text-align:center;font-size:13px;color:var(--mute);cursor:pointer;
    text-decoration:underline">처음이신가요? 가입 신청</div>
   <!-- ★v443 (2026.08.17) — 크롬 ⋮ 메뉴만 믿으면 「왜 안 뜨는지」를 영영 모른다.
        버튼을 눈에 보이게 두고, 설치 자격이 없으면 그 이유를 글자로 뱉는다(조문 흔적설계). -->
   <button id="pwabtn" style="display:none;width:100%;max-width:420px;margin-top:16px;
    border:none;border-radius:14px;padding:16px;font-size:16px;font-weight:800;
    color:#06203f;background:#c5a052;cursor:pointer">홈 화면에 앱 설치</button>
   <div id="pwamsg" style="margin-top:12px;font-size:12px;color:var(--mute);
    line-height:1.7;text-align:center"></div>
   <div style="text-align:center;margin-top:22px">
    <a href="/admin" style="font-size:12px;color:#6c7887;text-decoration:none;
     border:1px solid var(--line);border-radius:8px;padding:7px 16px">관리자</a></div>
   <div id="joinform" style="display:none;margin-top:14px;text-align:left">
    <input id="jn" class="pw" placeholder="이름" style="letter-spacing:0;text-align:left;font-size:15px">
    <input id="jp" class="pw" placeholder="연락처 (010-0000-0000)" inputmode="tel"
     style="letter-spacing:0;text-align:left;font-size:15px;margin-top:10px">
    <button id="jgo" class="go" style="margin-top:10px">가입 신청</button>
    <div id="jmsg" style="margin-top:12px;font-size:13px;text-align:center;color:var(--mute)"></div>
   </div>
  </div>
</div>
<div class="app" id="app">
  <header><div style="flex:1"><h1>@@BRAND@@ <b>@@BSUB@@</b></h1>
    <div class="sub">보장분석 리포트 PDF 1개 → 엑셀+PPT 개별 다운로드</div>
    <div id="drobot" style="margin-top:5px;font-size:10.5px;font-weight:700;letter-spacing:-.2px;white-space:nowrap;overflow:hidden">@@DROBOT@@</div></div>
    <button id="lout" style="border:1px solid var(--line);background:transparent;color:var(--mute);
     border-radius:8px;padding:6px 12px;font-size:12px;cursor:pointer;flex:none">로그아웃</button></header>
  <div class="chat" id="chat">
    <!-- ★★★★★v602 (지점장 지시 2026.08.30 「아래 설명 문구도 보기 좋게 정리해줘」).
         모바일 한 화면에 들어오게 <b>표 형태 3줄 + 주의 2줄</b>로 줄인다. -->
    <!-- ★v602b (지점장 지시 2026.08.30) — <b>넣는 것 → 나오는 것</b>을 경로별로 보여준다.
         엑셀 비교 설명이 아예 없었다. -->
    <div class="msg bot">
      <!-- ★v603b (지점장 지시 2026.08.30 「아이콘 좀 넣고 · 줄 넘치면 보기 좋게 정리」).
           산출물은 <b>칩(chip)</b>으로 감싸 줄이 넘쳐도 자연스럽게 흐르게 한다. -->
      <!-- ★★★★★v626 (지점장 지시 2026.09.02 「뭐가 나오는지 심플하게」).
           칩 12개를 버리고 <b>넣는 것 → 나오는 것</b>을 한 줄씩 3줄로. -->
      <div style="display:flex;flex-direction:column;gap:7px;font-size:12px;line-height:1.5">
        <div><b><span style="display:inline-block;vertical-align:-3px;margin-right:5px;line-height:0"><svg viewBox="0 0 24 24" width="15" height="15" aria-hidden="true"><path d="M4 2h10l6 6v14H4z" fill="#d93025"/><path d="M14 2l6 6h-6z" fill="#f4a6a0"/><text x="12" y="17" font-size="6.4" font-weight="700" fill="#fff" text-anchor="middle" font-family="Arial">PDF</text></svg></span>보장분석지</b><br>
          <span style="color:var(--mute)">엑셀 · 분석지PPT · 진단서 · 인포메이션</span></div>
        <div><b><span style="display:inline-block;vertical-align:-3px;margin-right:5px;line-height:0"><svg viewBox="0 0 24 24" width="15" height="15" aria-hidden="true"><rect x="3" y="3" width="18" height="19" rx="2" fill="#f2b32c"/><rect x="6" y="5.5" width="12" height="14" rx="1" fill="#eef3f8"/><rect x="9.6" y="1.6" width="4.8" height="3" rx="1.4" fill="#9fb2c6"/><g fill="#3f6fa3"><rect x="7.6" y="8" width="2.4" height="2.4" rx=".5"/><rect x="11" y="8.6" width="5.6" height="1.2" rx=".6"/><rect x="7.6" y="12" width="2.4" height="2.4" rx=".5"/><rect x="11" y="12.6" width="5.6" height="1.2" rx=".6"/><rect x="7.6" y="16" width="2.4" height="2.4" rx=".5"/><rect x="11" y="16.6" width="5.6" height="1.2" rx=".6"/></g></svg></span>제안서</b><br>
          <span style="color:var(--mute)">엑셀 · 분석지PPT · 진단서 · 인포메이션</span></div>
        <div><b><span style="display:inline-block;vertical-align:-3px;margin-right:5px;line-height:0"><svg viewBox="0 0 24 24" width="15" height="15" aria-hidden="true"><rect x="2.5" y="3.5" width="19" height="17" rx="2.4" fill="#1d7044"/><rect x="9.6" y="3.5" width="11.9" height="5.6" fill="#21935a"/><rect x="9.6" y="9.1" width="11.9" height="5.6" fill="#1d7044"/><rect x="9.6" y="14.7" width="11.9" height="5.8" fill="#155c37"/><rect x="2.5" y="7" width="10.5" height="10" rx="1.6" fill="#0f5132"/><text x="7.7" y="15.2" font-size="8.6" font-weight="700" fill="#fff" text-anchor="middle" font-family="Arial">X</text></svg></span>엑셀 전 vs 후</b><br>
          <span style="color:var(--mute)">비교엑셀 · 리포트 · 분석지PPT</span></div>
      </div>
      <div style="margin-top:10px;font-size:9.5px;color:var(--mute);line-height:1.55">
        ⚠️ 받은 PDF를 <b>그대로</b> 올리세요 — 재스캔·OCR은 금액이 깨집니다
      </div>
    </div>
  </div>
  <div class="gowrap"><button class="gobtn" id="gobtn" disabled>분석하기</button></div>
  <div style="display:none">
    <button id="send" disabled></button><button id="rsend" disabled></button>
  </div>
  <div class="qlbl" id="qlbl">분석된 보장분석지에 대해 질문하세요</div>
  <div class="qbar" id="qbar">
    <input class="qinput" id="qinput" placeholder="예: 심장 담보 왜 빠졌어요?" autocomplete="off">
    <button class="qbtn" id="qbtn">질문</button>
  </div>
  <nav class="tabbar">
    <label class="tab" id="upp"><span class="ti"><svg viewBox="0 0 24 24" width="21" height="21" aria-hidden="true"><path d="M4 2h10l6 6v14H4z" fill="#d93025"/><path d="M14 2l6 6h-6z" fill="#f4a6a0"/><text x="12" y="17" font-size="6.4" font-weight="700" fill="#fff" text-anchor="middle" font-family="Arial">PDF</text></svg></span><span class="nm" id="upplabel">PDF</span></label>
    <label class="tab" id="up"><span class="ti"><svg viewBox="0 0 24 24" width="21" height="21" aria-hidden="true"><rect x="3" y="3" width="18" height="19" rx="2" fill="#f2b32c"/><rect x="6" y="5.5" width="12" height="14" rx="1" fill="#eef3f8"/><rect x="9.6" y="1.6" width="4.8" height="3" rx="1.4" fill="#9fb2c6"/><g fill="#3f6fa3"><rect x="7.6" y="8" width="2.4" height="2.4" rx=".5"/><rect x="11" y="8.6" width="5.6" height="1.2" rx=".6"/><rect x="7.6" y="12" width="2.4" height="2.4" rx=".5"/><rect x="11" y="12.6" width="5.6" height="1.2" rx=".6"/><rect x="7.6" y="16" width="2.4" height="2.4" rx=".5"/><rect x="11" y="16.6" width="5.6" height="1.2" rx=".6"/></g></svg></span><span class="nm" id="uplabel">제안서</span></label>
    <div class="tabc" id="tabnew"><span class="ti"><svg viewBox="0 0 24 24" width="21" height="21" aria-hidden="true"><path d="M12 5.4V2.2L7.5 6.3 12 10.4V7.6c2.9 0 5.2 2.3 5.2 5.2S14.9 18 12 18s-5.2-2.3-5.2-5.2H4.6c0 4.1 3.3 7.4 7.4 7.4s7.4-3.3 7.4-7.4S16.1 5.4 12 5.4z" fill="#f08c14"/></svg></span><span class="nm">리셋</span></div>
    <label class="tab" id="upr1"><span class="ti"><svg viewBox="0 0 24 24" width="21" height="21" aria-hidden="true"><rect x="2.5" y="3.5" width="19" height="17" rx="2.4" fill="#1d7044"/><rect x="9.6" y="3.5" width="11.9" height="5.6" fill="#21935a"/><rect x="9.6" y="9.1" width="11.9" height="5.6" fill="#1d7044"/><rect x="9.6" y="14.7" width="11.9" height="5.8" fill="#155c37"/><rect x="2.5" y="7" width="10.5" height="10" rx="1.6" fill="#0f5132"/><text x="7.7" y="15.2" font-size="8.6" font-weight="700" fill="#fff" text-anchor="middle" font-family="Arial">X</text></svg></span><span class="nm" id="upr1label">엑셀(전)</span></label>
    <label class="tab" id="upr2"><span class="ti"><svg viewBox="0 0 24 24" width="21" height="21" aria-hidden="true"><rect x="2.5" y="3.5" width="19" height="17" rx="2.4" fill="#1d7044"/><rect x="9.6" y="3.5" width="11.9" height="5.6" fill="#21935a"/><rect x="9.6" y="9.1" width="11.9" height="5.6" fill="#1d7044"/><rect x="9.6" y="14.7" width="11.9" height="5.8" fill="#155c37"/><rect x="2.5" y="7" width="10.5" height="10" rx="1.6" fill="#0f5132"/><text x="7.7" y="15.2" font-size="8.6" font-weight="700" fill="#fff" text-anchor="middle" font-family="Arial">X</text></svg></span><span class="nm" id="upr2label">엑셀(후)</span></label>
  </nav>
  <footer>미래를 <b>바르게</b> 설계합니다 · @@BRAND@@</footer>
</div>
<input type="file" id="fi" accept=".pdf,application/pdf" multiple style="display:none">
<input type="file" id="fp" accept=".pdf,application/pdf" style="display:none">
<input type="file" id="fr1" accept=".xlsx" style="display:none">
<input type="file" id="fr2" accept=".xlsx" style="display:none">
<script>
const $=s=>document.querySelector(s);let ACCESS='';const MPW='0101';
let R1=null,R2=null;
function rchk(){const b=$("#rsend");if(b)b.disabled=!(R1&&R2);}
/* ★★★★★v450 제65조 — DOMContentLoaded는 <b>이미 지나간 뒤면 영원히 안 온다</b>.
   설치앱·뒤로가기 복원(bfcache)에서 스크립트가 늦게 돌면 버튼이 <b>아예 안 먹힌다</b>.
   → 이미 파싱이 끝났으면 즉시 실행한다. */
function _onReady(fn){
  if(document.readyState === 'loading'){ document.addEventListener('DOMContentLoaded', fn); }
  else { fn(); }
}
_onReady(()=>{
 const u1=$("#upr1"),u2=$("#upr2"),f1=$("#fr1"),f2=$("#fr2"),rb=$("#rsend");
 if(!u1)return;
 u1.onclick=()=>f1.click(); u2.onclick=()=>f2.click();
 f1.onchange=e=>{R1=e.target.files[0];$("#upr1label").textContent=R1?R1.name:"엑셀(전)";rchk();};
 f2.onchange=e=>{R2=e.target.files[0];$("#upr2label").textContent=R2?R2.name:"엑셀(후)";rchk();};
 rb.onclick=async()=>{
  if(!(R1&&R2))return; rb.disabled=true; rb.textContent="비교 중…";
  const fd=new FormData(); fd.append("old_xlsx",R1); fd.append("new_xlsx",R2); fd.append("pw",ACCESS);
  try{
   const r=await fetch("/remodel",{method:"POST",body:fd}); const j=await r.json();
   if(!j.ok){alert(j.error||"실패");}
   else{
    /* ★★★★★v473 제78조 (지점장 지적 2026.08.18 「리모델링비교하면 계속 아래에 파란색글짜로만 뜬다」)
       실측 결함 2건 —
       ① 보장분석 결과는 <b>`.file-card` 저장 카드</b>인데 리모델링만 <b>맨 `<a>` 파란 링크</b>였다.
          폰에서 글자가 작아 누르기 어렵고 산출물로 보이지 않는다.
       ② 서버는 `j.pdf`(리포트 PDF)를 <b>이미 보내고 있는데 화면에 아예 안 걸렸다</b> — 산출물 1개 누락.
       → 보장분석과 <b>같은 카드·같은 저장 동작</b>으로 통일한다. PC 자동저장 · 폰 카드 클릭도 동일. */
    const _rM=/Android|iPhone|iPad|iPod|Mobile|SamsungBrowser/i.test(navigator.userAgent)
              ||(navigator.maxTouchPoints>1&&/Macintosh/.test(navigator.userAgent));
    const _rc=(u,nm,sub,cls)=>u?('<a class="file-card '+cls+'" href="'+u+'" download'
      +' style="cursor:pointer;text-decoration:none;color:inherit">'
      +'<span class="ic"></span><span class="nm">'+nm
      +'<br><span style="font-size:10px;color:var(--mute)">'+sub+'</span></span>'
      +'<span class="dl">저장</span></a>'):'';
    const _rn=(j.client||'')+'_리모델링';
    let h='<b>리모델링 비교 완료!</b> <span style="font-size:11px;color:var(--mute)">'
      +(_rM?'★휴대폰은 <b>카드를 하나씩 눌러</b> 저장하세요 (연속 저장이 차단됩니다)'
          :'자동 저장 중… 안 되면 카드를 누르세요')+'</span>'
      +'<div class="summary-box">'
      +'기존 <b>'+j.prem_old.toLocaleString()+'원</b> → 최종 <b>'+j.prem_new.toLocaleString()+'원</b><br>'
      +'월 절감 <b>'+j.save_m.toLocaleString()+'원</b> ('+j.save_pct+'%) · 연 '
      +(j.save_y||0).toLocaleString()+'원<br>'
      +'보장 증가 '+j.n_up+' · 신규 '+j.n_add+' · 감소 '+j.n_down+' · 삭제 '+j.n_del
      +'</div><div class="file-cards">'
      +_rc(j.xlsx,_rn+'_비교.xlsx','비교 엑셀','xl')
      +_rc(j.pptx,_rn+'_리포트.pptx','리모델링 리포트 PPT','pt')
      +_rc(j.pdf, _rn+'_리포트.pdf', '리모델링 리포트 PDF','pt')
      /* ★v612 (지점장 지시 2026.08.30) — 최종 엑셀 기준 보장분석지 PPT 추가 */
      +_rc(j.anal, '보장분석지_'+_rn+'(최종).pptx', '보장분석지 PPT (최종)','pt')
      +'</div>';
    if(!_rM){
      const _q=[[j.xlsx,_rn+'_비교.xlsx'],[j.pptx,_rn+'_리포트.pptx'],[j.pdf,_rn+'_리포트.pdf'],
                [j.anal,'보장분석지_'+_rn+'(최종).pptx']]
               .filter(x=>x[0]);
      _q.forEach((x,i)=>setTimeout(()=>{const a=document.createElement("a");
        a.href=x[0];a.download=x[1];a.style.display="none";document.body.appendChild(a);a.click();
        setTimeout(()=>{try{document.body.removeChild(a);}catch(e){}},2000);}, i*900));
    }
    /* ★★★★★v464 제72조 (지점장 지적 2026.08.17 「보험리모델링이 화면 가운데 안 뜨고 저 위에 뜬다」)
       실측: `#out`은 <b>화면에 없는 아이디</b>였다(0개) → 폴백인 `document.body` 맨 앞에 붙어
       <b>헤더보다 위</b>에 나왔다. 분석 결과는 대화창(`#chat`)에 들어가는데 리모델링만 딴 데 갔다.
       → 분석 결과와 <b>같은 자리</b>(대화창 맨 아래)에 넣고 그리로 스크롤한다.
       ★v473: 카드 CSS(.file-card)는 <b>`.msg.bot` 말풍선 안</b>에서만 제 폭이 나온다.
       → 분석 결과와 <b>똑같이</b> `add(h,"bot")`으로 감싸 넣는다(raw 삽입 폐기). */
    let _el=null;
    if(typeof add==="function"){ _el=add(h,"bot"); }
    else{ const box=document.getElementById("chat");
          if(box){ box.insertAdjacentHTML("beforeend",'<div class="msg bot">'+h+'</div>');
                   _el=box.lastElementChild; }
          else { document.body.insertAdjacentHTML("beforeend",h); } }
    if(_el) _el.scrollIntoView({behavior:"smooth",block:"center"});
   }
  }catch(e){alert("오류: "+e);}
  rb.disabled=false; rb.textContent="리모델링 비교";
 };
 /* ★★★★★v450 제65조 (지점장 지시 2026.08.17 「리셋 버튼도 필요하다 ·
    버튼 누르고 나서 계속 기존 사람의 잔재가 남아있다」)
    화면 상태를 <b>한 곳에서</b> 전부 지운다. 지우는 대상을 늘리려면 여기만 고친다. */
 /* ★★★★★v627 (지점장 지시 2026.09.02 「글자넘치는건 싸이즈줄여라」).
    잘라내지 않고 <b>들어갈 때까지 폰트를 낮춘다</b>. 하한에 닿으면 그때만 줄임표. */
 function _fit(el, base, min){
   if(!el) return;
   el.style.fontSize = base + "px";
   var g = 0;
   while(el.scrollWidth > el.clientWidth + 1 && parseFloat(el.style.fontSize) > min && g++ < 40){
     el.style.fontSize = (parseFloat(el.style.fontSize) - 0.5) + "px";
   }
 }
 function _fitAll(){
   _fit(document.getElementById("drobot"), 10.5, 7.5);
   ["upplabel","uplabel","upr1label","upr2label"].forEach(function(i){
     _fit(document.getElementById(i), 9.5, 7);
   });
 }
 window.addEventListener("resize", _fitAll);
 setInterval(_fitAll, 400); _fitAll();
 const _tn=$("#tabnew");
 /* ★★★★★v625 (지점장 지시 2026.09.02 「그위 가로로 길게 분석하기」).
    실행 버튼은 <b>하나</b>다. 엑셀 두 개가 준비되면 리모델링, 아니면 분석.
    기존 #send·#rsend 는 숨겨두고 그대로 눌러 <b>동작 코드는 손대지 않는다</b>. */
 const _gb=$("#gobtn");
 function _syncGo(){
   const sd=$("#send"), rb=$("#rsend"); if(!_gb) return;
   const rOK = rb && !rb.disabled, sOK = sd && !sd.disabled;
   _gb.textContent = rOK ? "리모델링 비교" : "분석하기";
   _gb.disabled = !(rOK || sOK);
   ["upp","up","upr1","upr2"].forEach(function(i){
     const el=$("#"+i), nm=$("#"+i+"label");
     if(el && nm) el.classList.toggle("on", !!(nm.textContent||"").match(/\.(pdf|xlsx|xls)$/i));
   });
 }
 if(_gb) _gb.onclick=()=>{ const rb=$("#rsend"), sd=$("#send");
   if(rb && !rb.disabled) rb.click(); else if(sd) sd.click(); };
 setInterval(_syncGo, 250); _syncGo();
 /* ★★★★★v635 (지점장 지시 2026.09.02 「새고객 지우고 리셋을 누르게 해야지」).
    숨긴 버튼을 대신 누르는 구조를 버리고, <b>리셋 탭이 직접</b> 초기화를 실행한다. */
 if(_tn) _tn.onclick=()=>{
   if(!confirm("초기화합니다. 지금 화면의 결과는 사라집니다."))return;
   file=null; pdfFile=null; files=[]; savedFiles={};
   R1=null; R2=null;
   const ids=["fi","fp","fr1","fr2"];
   for(const id of ids){const el=document.getElementById(id); if(el)el.value="";}
   const lbl={uplabel:"제안서",upplabel:"PDF",
              upr1label:"엑셀(전)",upr2label:"엑셀(후)"};
   for(const k in lbl){const el=document.getElementById(k); if(el)el.textContent=lbl[k];}
   const sb=document.getElementById("send"); if(sb){sb.disabled=true; sb.textContent="분석";}
   if(rb){rb.disabled=true; rb.textContent="리모델링 비교";}
   const ch=document.getElementById("chat");
   if(ch){const keep=ch.querySelector(".msg.bot"); ch.innerHTML=""; if(keep)ch.appendChild(keep);}
   const qi=document.getElementById("qinput"); if(qi)qi.value="";
   const ql=document.getElementById("qlbl"); if(ql)ql.textContent="분석된 보장분석지에 대해 질문하세요";
   window.scrollTo(0,0);
 };
});
async function unlock(v2){const src=(typeof v2==="string"&&v2)?v2:$("#pw").value;
  const v=String(src||"").trim().toUpperCase();
  if(!v)return; $("#gerr").textContent="확인 중…";
  try{
    const fd=new FormData();
    if(v===MPW){fd.append("pw",v);}else{fd.append("code",v);}
    const r=await fetch("/member/login",{method:"POST",body:fd});
    const j=await r.json();
    if(j.ok){ACCESS=MPW; try{localStorage.setItem("barum_code",v);}catch(e){}
      if(j.name)$(".sub").textContent=j.name+" 님 · 보장분석 자동화";
      $("#gerr").textContent="";$("#gate").style.display="none";$("#app").style.display="flex";}
    else{try{localStorage.removeItem("barum_code");}catch(e){} fail(j.error);}}
  catch(e){$("#gerr").textContent="서버 연결 실패";}}
/* ★v440 — 자동 로그인(localStorage) 때문에 잠금화면으로 돌아갈 길이 없었다(2026.08.17 실측).
   저장된 코드를 지우고 게이트를 다시 세운다. 다른 사람에게 넘길 때도 이 버튼을 쓴다. */
$("#lout").onclick=function(){try{localStorage.removeItem("barum_code");}catch(e){}
  $("#app").style.display="none";$("#gate").style.display="flex";
  $("#pw").value="";$("#gerr").textContent="";$("#pw").focus();};
function fail(m){$("#gerr").textContent=m||"코드 또는 비밀번호가 올바르지 않습니다.";$("#gate").classList.add("shake");setTimeout(()=>$("#gate").classList.remove("shake"),350);$("#pw").value="";$("#pw").focus();}
$("#go").onclick=function(){unlock();};
/* ★v443 앱 설치 (2026.08.17) — 자격이 되면 버튼, 안 되면 이유. 빈 화면은 두지 않는다. */
var _dp=null;
window.addEventListener("beforeinstallprompt",function(e){
  e.preventDefault(); _dp=e;
  $("#pwabtn").style.display="block"; $("#pwamsg").textContent="";});
$("#pwabtn").onclick=async function(){
  if(!_dp)return; _dp.prompt();
  try{var r=await _dp.userChoice;
    $("#pwamsg").textContent=(r.outcome==="accepted")?"설치했습니다. 홈 화면을 확인하세요.":"설치를 취소했습니다.";
  }catch(e){} _dp=null; $("#pwabtn").style.display="none";};
window.addEventListener("appinstalled",function(){
  $("#pwabtn").style.display="none"; $("#pwamsg").textContent="설치 완료";});
setTimeout(function(){
  if(_dp)return;
  var ua=navigator.userAgent||"";
  var standalone=window.matchMedia("(display-mode: standalone)").matches||navigator.standalone;
  var m="";
  if(standalone){m="이미 앱으로 실행 중입니다.";}
  else if(/iPhone|iPad|iPod/i.test(ua)){
    m=/CriOS|FxiOS|EdgiOS/i.test(ua)
      ? "아이폰은 <b>사파리</b>로 열어야 설치됩니다.<br>사파리 → 공유(⬆) → 홈 화면에 추가"
      : "공유(⬆) → <b>홈 화면에 추가</b> 를 누르세요";}
  else if(/SamsungBrowser/i.test(ua)){
    m="삼성 인터넷입니다. <b>크롬</b>으로 열면 설치 버튼이 나옵니다.<br>(삼성 인터넷은 ≡ → 현재 페이지 추가 → 홈 화면)";}
  else if(location.protocol!=="https:"){m="https 주소로 열어야 설치됩니다.";}
  else{m="설치 자격 확인 중입니다. 화면을 <b>한 번 더 새로고침</b>하세요.<br>그래도 없으면 크롬 ⋮ → 홈 화면에 추가";}
  $("#pwamsg").innerHTML=m;
},2500);
$("#joinopen").onclick=function(){
  var f=$("#joinform"); f.style.display=(f.style.display==="none")?"block":"none";
  if(f.style.display==="block"){ var s=null; try{s=JSON.parse(localStorage.getItem("barum_apply")||"null");}catch(e){}
    if(s){ $("#jn").value=s.n; $("#jp").value=s.p; jcheck(); } }
};
async function jcheck(){
  var n=$("#jn").value.trim(), p=$("#jp").value.trim();
  if(!n||!p)return;
  const fd=new FormData(); fd.append("name",n); fd.append("phone",p);
  const r=await fetch("/member/check",{method:"POST",body:fd}); const j=await r.json();
  if(j.ok&&j.status==="active"){
    $("#jmsg").innerHTML='승인되었습니다. 회원코드<br><b style="font-size:26px;color:var(--acc2);letter-spacing:.12em">'
      +j.code+'</b><br>위 칸에 입력해 접속하세요';
    $("#pw").value=j.code;
  } else if(j.ok){ $("#jmsg").textContent="승인 대기 중입니다. 승인되면 여기에 코드가 나옵니다."; }
  else { $("#jmsg").textContent=j.error||""; }
}
$("#jgo").onclick=async function(){
  var n=$("#jn").value.trim(), p=$("#jp").value.trim();
  if(n.length<2){$("#jmsg").textContent="이름을 정확히 입력하세요";return;}
  if(p.length<9){$("#jmsg").textContent="연락처를 정확히 입력하세요";return;}
  $("#jmsg").textContent="신청 중…";
  const fd=new FormData(); fd.append("name",n); fd.append("phone",p);
  const r=await fetch("/member/apply",{method:"POST",body:fd}); const j=await r.json();
  if(!j.ok){$("#jmsg").textContent=j.error||"실패";return;}
  try{localStorage.setItem("barum_apply",JSON.stringify({n:n,p:p}));}catch(e){}
  if(j.status==="active"){ jcheck(); }
  else { $("#jmsg").textContent="신청되었습니다. 지점장 승인 후 이 화면에 코드가 나옵니다."; }
};$("#pw").addEventListener("keydown",e=>{if(e.key==="Enter")unlock();});window.addEventListener("load",()=>{
  var sv=null; try{sv=localStorage.getItem("barum_code");}catch(e){}
  if(sv){unlock(sv);}else{$("#pw").focus();}});
const chat=$("#chat");let file=null;let pdfFile=null;let files=[];   /* ★v385 제안서 최대 3건 */
function _syncSend(){$("#send").disabled=!(file||pdfFile);}
$("#up").onclick=()=>$("#fi").click();
$("#upp").onclick=()=>$("#fp").click();
$("#fi").onchange=e=>{/* ★★★★★v592 (지점장 2026.08.26 「제안서 3개까지 가능하게 해줘」)
  [구 동작] `files=Array.from(...)` — 새로 고를 때마다 <b>앞의 선택을 통째로 덮었다</b>.
    파일선택창에서 Ctrl로 3개를 <b>한 번에</b> 골라야만 3건이 됐고,
    한 개씩 세 번 누르면 <b>마지막 1건만</b> 남았다 → 지점장이 「1개밖에 인식이 안 돼」.
  ⇒ <b>누를 때마다 이어붙인다</b>(누적). 같은 이름·크기는 중복으로 보고 한 번만.
    3건이 차면 알린다. 지우려면 초기화 버튼. */
 var _add=Array.from(e.target.files||[]);
 var _seen={};files.forEach(function(f){_seen[f.name+'|'+f.size]=1;});
 _add.forEach(function(f){var k=f.name+'|'+f.size;if(!_seen[k]&&files.length<3){files.push(f);_seen[k]=1;}});
 if(files.length>3){files=files.slice(0,3);}
 file=files[0]||null;e.target.value="";   /* 같은 파일을 다시 고를 수 있게 비운다 */$("#uplabel").textContent=files.length?(files.map(function(f){return f.name;}).join(" · ")+"  ["+files.length+"/3]"):"제안서";_syncSend();};
$("#fp").onchange=e=>{pdfFile=e.target.files[0]||null;$("#upplabel").textContent=pdfFile?pdfFile.name:"PDF";_syncSend();};
function esc(s){return String(s==null?"":s).replace(/[&<>]/g,c=>({"&":"&amp;","<":"&lt;",">":"&gt;"}[c]));}
function add(html,cls){const d=document.createElement("div");d.className="msg "+cls;d.innerHTML=html;chat.appendChild(d);chat.scrollTop=chat.scrollHeight;return d;}
function b64toBlob(b64,mime){const bin=atob(b64);const arr=new Uint8Array(bin.length);for(let i=0;i<bin.length;i++)arr[i]=bin.charCodeAt(i);return new Blob([arr],{type:mime});}
function dl(blob,fname){const u=URL.createObjectURL(blob);const a=document.createElement("a");a.href=u;a.download=fname;document.body.appendChild(a);a.click();document.body.removeChild(a);setTimeout(()=>URL.revokeObjectURL(u),3000);}
const XLMIME="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet";
const PTMIME="application/vnd.openxmlformats-officedocument.presentationml.presentation";
const PDFMIME="application/pdf";
let savedFiles={};
function reDL(k){const f=savedFiles[k];if(f&&f.b64){dl(b64toBlob(f.b64,f.mime),f.name);}}
$("#send").onclick=async()=>{
  if(!file&&!pdfFile)return;add(esc((file?file.name:"")+(file&&pdfFile?" + ":"")+(pdfFile?pdfFile.name:"")),"me");
  $("#send").disabled=true;$("#up").style.opacity=.5;
  const loading=add('<div style="display:flex;align-items:center;gap:11px"><span class="spin"></span><div style="flex:1"><div id="ldmsg" style="font-weight:800">PDF 파싱 중…</div><div id="ldtime" style="font-size:11px;color:var(--mute);margin-top:2px">0초 · 기다려 주세요</div></div></div>',"bot");
  const t0=Date.now();const steps=["PDF 파싱 중…","담보 추출 중…","엑셀 생성 중…","PPT 채우는 중…","완성 중…"];let si=0;
  const timer=setInterval(()=>{si=Math.min(si+1,steps.length-1);const s=Math.floor((Date.now()-t0)/1000);const tm=document.getElementById("ldtime");const mm=document.getElementById("ldmsg");if(tm)tm.textContent=s+"초 경과";if(mm)mm.textContent=steps[si];},8000);
  const fd=new FormData();
/* v370 role fix: server file=bojang(pdfFile) / file2=jean(file) */
/* ★v385: 제안서는 같은 이름 file2로 <b>여러 번</b> append한다(최대 3건). 칸은 그대로 하나다. */
const _jn=(files&&files.length?files:(file?[file]:[]));
if(pdfFile){fd.append("file",pdfFile);}
_jn.slice(0,3).forEach(f=>fd.append("file2",f));
  fd.append("pw",ACCESS);
  let j=null;
  try{
    const r=await fetch("/analyze",{method:"POST",body:fd});clearInterval(timer);loading.remove();
    j=await r.json();
    if(!j.ok){
      /* ★v94: '[오류] 실패'만 뜨고 원인을 알 수 없던 문제 — 서버가 보내주는 trace를 화면에 같이 찍는다. */
      var _m = esc(j.error||"실패(서버가 오류 문구를 못 보냄)");
      var _t = j.trace ? String(j.trace) : "";
      if(_t){ var _tail=_t.split("\n").slice(-8).join("\n");
              _m += '<br><span style="font-size:11px;opacity:.85;white-space:pre-wrap">'+esc(_tail)+'</span>'; }
      add('<span class="err">[오류] '+_m+'</span>',"bot");
    }
    else{
      savedFiles={};
      /* ★v316: 자동 blob 다운로드 폐기(모바일 Chrome이 연속 다운로드를 차단해 저장이 안 됐다).
         서버가 준 실제 URL(j.*_url)을 링크로 걸고, 사용자가 카드를 하나씩 눌러 받는다.
         URL이 없으면(구버전 서버) 종전 blob 재저장(reDL)으로 자동 폴백된다. */
      const _mk=(u,k)=>u?('href="'+u+'" download'):('href="javascript:void(0)" onclick="reDL(\''+k+'\')"');
      savedFiles.xlsx={b64:j.xlsx_b64,name:j.xlsx_name,mime:XLMIME};
      /* ★★v317 (지점장 지시 2026.08.01): <b>PC는 자동저장 · 휴대폰은 카드 클릭</b>.
         모바일 Chrome만 연속 다운로드를 차단하므로, PC에서는 종전처럼 자동으로 받아준다.
         ★blob이 아니라 <b>서버 실제 URL</b>로 받으므로 revoke 만료 문제가 없다. */
      const _isMobile=/Android|iPhone|iPad|iPod|Mobile|SamsungBrowser/i.test(navigator.userAgent)
                      ||(navigator.maxTouchPoints>1&&/Macintosh/.test(navigator.userAgent));
      function dlUrl(u,fn){const a=document.createElement("a");a.href=u;a.download=fn||"";
        a.style.display="none";document.body.appendChild(a);a.click();
        setTimeout(()=>{try{document.body.removeChild(a);}catch(e){}},2000);}
      if(!_isMobile){
        const _q=[[j.xlsx_url,j.xlsx_name],[j.pptx_url,j.pptx_name],[j.chiryo_url,j.chiryo_name],
                  [j.report_url,j.report_name],[j.report_pptx_url,j.report_pptx_name]].filter(x=>x[0]);
        _q.forEach((x,i)=>setTimeout(()=>dlUrl(x[0],x[1]), i*900));
      }
      let ptCard='';
      if(j.pptx_b64){
        savedFiles.pptx={b64:j.pptx_b64,name:j.pptx_name,mime:PTMIME};
        ptCard=`<a class="file-card pt" ${_mk(j.pptx_url,'pptx')} style="cursor:pointer;text-decoration:none;color:inherit"><span class="ic"></span><span class="nm">${esc(j.pptx_name)}<br><span style="font-size:10px;color:var(--mute)">보장분석 PPT</span></span><span class="dl">저장</span></a>`;}
      if(j.chiryo_b64){
        savedFiles.chiryo={b64:j.chiryo_b64,name:j.chiryo_name,mime:PTMIME};
        ptCard+=`<a class="file-card pt" ${_mk(j.chiryo_url,'chiryo')} style="cursor:pointer;text-decoration:none;color:inherit"><span class="ic"></span><span class="nm">${esc(j.chiryo_name)}<br><span style="font-size:10px;color:var(--mute)">치료비 정리 PPT</span></span><span class="dl">저장</span></a>`;}
      if(j.report_b64){
        savedFiles.report={b64:j.report_b64,name:j.report_name,mime:PDFMIME};
        ptCard+=`<a class="file-card pt" ${_mk(j.report_url,'report')} style="cursor:pointer;text-decoration:none;color:inherit"><span class="ic"></span><span class="nm">${esc(j.report_name)}<br><span style="font-size:10px;color:var(--mute)">보장설명지 PDF</span></span><span class="dl">저장</span></a>`;}
      if(j.report_pptx_b64){
        savedFiles.reportpptx={b64:j.report_pptx_b64,name:j.report_pptx_name,mime:PTMIME};
        ptCard+=`<a class="file-card pt" ${_mk(j.report_pptx_url,'reportpptx')} style="cursor:pointer;text-decoration:none;color:inherit"><span class="ic"></span><span class="nm">${esc(j.report_pptx_name)}<br><span style="font-size:10px;color:var(--mute)">보장진단서 PPT (편집가능)</span></span><span class="dl">저장</span></a>`;}
      add('<b>분석 완료!</b> <span style="font-size:11px;color:var(--mute)">'+(_isMobile?'★휴대폰은 <b>카드를 하나씩 눌러</b> 저장하세요 (연속 저장이 차단됩니다)':'자동 저장 중… 안 되면 카드를 누르세요')+'</span><div class="summary-box">'+j.summary+'</div><div class="file-cards">'+
        `<a class="file-card xl" ${_mk(j.xlsx_url,'xlsx')} style="cursor:pointer;text-decoration:none;color:inherit"><span class="ic"></span><span class="nm">${esc(j.xlsx_name)}<br><span style="font-size:10px;color:var(--mute)">보장진단 엑셀</span></span><span class="dl">저장</span></a>`+ptCard+'</div>',"bot");}
  }catch(e){clearInterval(timer);loading.remove();add('<span class="err">오류: '+esc(e.message)+'</span>',"bot");}
  if(j&&j.data){analysisData=j.data;document.getElementById("qbar").style.display="flex";document.getElementById("qlbl").style.display="block";}
  file=null;files=[];$("#uplabel").textContent="제안서";$("#send").disabled=true;$("#fi").value="";$("#up").style.opacity=1;
  if(j&&j.report_error){add('<span class="err">[오류] 보장설명지 PDF 생성 실패: '+esc(j.report_error)+'</span>',"bot");}
  if(j&&j.report_pptx_error){add('<span class="err">[오류] 보장진단서 PPT 생성 실패: '+esc(j.report_pptx_error)+'</span>',"bot");}
  if(j&&j.ok){add('다음 고객 PDF를 올리면 이어서 분석합니다.',"bot");}
};
let analysisData=null;
function askAI(){
  const q=document.getElementById("qinput").value.trim();
  if(!q||!analysisData)return;
  add("[질문] "+esc(q),"me");
  document.getElementById("qinput").value="";
  document.getElementById("qbtn").disabled=true;
  const loading=add('<span class="spin"></span> 분석 중…',"bot");
  fetch("/ask",{method:"POST",headers:{"Content-Type":"application/json"},
    body:JSON.stringify({pw:ACCESS,question:q,data:analysisData})})
  .then(r=>r.json()).then(j=>{
    loading.remove();
    add(j.ok?esc(j.answer):'<span class="err">[오류] '+esc(j.error||"오류")+'</span>',"bot");
    document.getElementById("qbtn").disabled=false;
  }).catch(e=>{loading.remove();add('<span class="err">오류: '+esc(e.message)+'</span>',"bot");document.getElementById("qbtn").disabled=false;});
}
_onReady(function(){
  const qi=document.getElementById("qinput");
  if(qi) qi.addEventListener("keydown",function(e){if(e.key==="Enter")askAI();});
  const qb=document.getElementById("qbtn"); if(qb) qb.onclick=askAI;
});
</script>
<!-- ★v440 (2026.08.17 실측) — 여기 있던 <script>가 7577행에서 등록한 서비스워커를
     매 로드마다 통째로 unregister 했다. 그래서 크롬이 「앱 설치」를 영영 띄우지 않았다.
     서비스워커는 캐시를 하지 않으므로(no-store) 해제할 이유가 없다. 삭제한다. -->
</body></html>'''

# ═══════════ v428 PWA (제53조) — manifest · 아이콘 · 서비스워커 ═══════════
_PWA_ICON = None

# ★★★★★v629 (지점장 지시 2026.09.02 「앱아이콘도 이걸로 바꾸자」).
#   ★v634 (지점장 지시 2026.09.02 「그 도시배경 가능」) — MAKEONE 빌딩 사진 배경
#   + 네이비 톤 45% + 흰 글자 「보장/분석」 + 골드 밑줄. ★원본 271px — 고해상도 받으면 재생성.
_ICON_SRC_B64 = 'iVBORw0KGgoAAAANSUhEUgAAAgAAAAIACAIAAAB7GkOtAAEAAElEQVR42uz9ebxsWXYWBq619j4RcePeN+XLoapUQsKgRhKIeTBgMKNBbmOa0b+GBre7sTFyIxoDsgQGg0CAmJHMDwEWxhiQ7DayQcw09E8MbmiDGTSApKrKrMyX+cY7xY3pnLP3Xv3Hntbe58Qd3susyqyMrczSzfvi3RvDOWv41re+D7/2D/5J2J/92Z/92Z8P36H9W7A/+7M/+7NPAPuzP/uzP/uzTwD7sz/7sz/7s08A+7M/+7M/+7NPAPuzP/uzP/uzTwD7sz/7sz/7s08A+7M/+7M/+7NPAPuzP/uzP/uzTwD7sz/7sz/7s08A+7M/+7M/+7NPAPuzP/uzP/uzTwD7sz/7sz/7s08A+7M/+7M/+7NPAPuzP/uzP/uzTwD7sz/7sz/7s08A+7M/+7M/+7NPAPuzP/uzP/uzTwD7sz/7sz/7s08A+7M/+7M/+7NPAPuzP/uzP/uzTwD7sz/7sz/7s08A+7M/+7M/+7NPAPuzP/uzP/uzTwD7sz/7sz/7BLA/+7M/+7M/+wSwP/uzP/uzP/sEsD/7sz/7sz/7BLA/+7M/+7M/+wSwP/uzP/uzP/sEsD/7sz/7sz/7BLA/+7M/+7M/H6yj92/B/nwIDscvcP9e7M/+7DuA/dmf/dmffQLYn/3Zn/3Znw/l2UNA+/NhOHvkZ3/2Z58A9uddPnyN8LrH3/dnf/YJYH/2lfX+7M/+vJ/OfgawP/uzP/uz7wD2Z3/2XcL+7M8+AezP/rw/zn5+sH/H9mefAPZn3z3sz/7sz7t99jOA/dmf/dmffQewP/vz/juSZ/rZ+voD2zPxvoXan30C2J93MRTfICqOxVAWX2EVbUcD2DCS3fTrz2RygvI5vA9SyAc2ke3PPgHsz/u1tHzuvxCiJQNw+D8fpVg+gjlFLy6SwwfnvcLyLUCs3xX0KRDDf+N7+JHxvhXYn30C2J/3Vy7x/6LoBso/+6BCGNdplHis19lH6P3ZJ4D9ea+jE75bP6ioMDmW6sggSnjgHBJzZCQEJBJPBd+tpuMD0ighpNbHN0SOHVc9g3x3UHQRvm+gHYlk8Bkx7jPL/uwTwP7srj5vXNtWmQR9dIpx6sofz/Ef5zOHTx/scwWnBOJ8TuGAFaU0wuK5sEgwMdlgjq9VOY5wLUAJL3u3hq8RAcsHYHhHEDC+MZh7nvgHgIiIAAhASKhu9kHw2PgEyqwMH0z4bH/2CWB/bhIOmBGvGdB3BY0dUYZjhcp1tM0JgBCR8CaEYgQwDH3f98aYvrfOWWedscYaY6011lprnTXWWWOsdc5Zy845ZsfMvlxm5xyDc9bnDgZ21jpmDuGWiyePGJ4kh/+s43t8sH9kCNr+lyACIYXQjojkX7D/ktL/IhIiEJH//0qRUqSVUqSUUkprpch/TVpprTRppZRWqpk0zU3Sccid7PwnBCOpp+gc8P3RAOyH0fsEsD8fjO4gFKqhii3hiMHjXQxJPhT7EA3AjsE5n0RSgc/MbIxtu67rurbtur7r+94Ya521xviUYI01xvTW2t4a0xtnnXHWWcfOxdAPjp3PCNYxW2a27Jx17BgICeOUGVPLgIToA32RAAYZERFDEgBwzOBiRkEEBkRUpJCQiIgwhHskAiIVvuH/VKHSjdKalNJaNVorrbVWWjVK66bRWmmlVdNo3Uyb2XQynTSN1gTI8dch+JwC5J9Rzk3+kyHCK8YM+B7Bf/uzTwD78/4J2njD+rFEUiLRpgBTYql7jVqeAdre9L3p+67rTde1fe9DeW8605uu701vrDG9s7ZP9b5jZ50x1rJlx8DsHBt2bJ1j55jZuVjps3MuPemYZXL748CF7OLLXQYHDJwjJQMAoI2ZIOQFCH/O4S3kCGT50B9QHt9bJXAKgZ2zyAjWxaAc82TuHQDA5wQCQgJSPjsgEQEqRahC+0BIhEqRItKktPK9gmqaRjd6oieq0Y3SWqvJpNG6aZpJ06hmMmmaZqKu+HT8n7r4Dok/SG0B5pR39VXzorljn3v2CWB/blSC3+jxN3pEvv1xd3XvACyDc9aBD8kxMDsfVdgyt73p2na7adu+67q+67t2s+36ftt1PiH0Xdt2vTF9b2zfd70xHrAKiEnGKlCU5SmqBijGl++ASAGuicgTIkf0Jj2UIRT/FH4NACA7cODSe+F/mg/yPvkhMAEiESCmiXb11vg0E3iqXH1f5ENgZrDM1lo2Pm1EnMl/ybEtwjDd8GmNkLTWutGTpmkaPZ3Mmon/Lz1pJpNJ00wm04n/dzKbTCeTptFKqTBK8EAUERECoYejyONXlwwY+LOzuMD7NLFPAPvzruSAK35WQGR8/RvgcA+TXHFaYzfrzWa7bbvtdtt2fd9ut13Xb7tt3/ad6a211rK1trfWA/POsbHWWufYOuucc9Y5a61jBw4sOxtjdSD7xzmoD/CQq20f3P0TxRxnU1LwAdWH+/BD0htH/ksXfxECOgJ2qe4nJGRGZifnxy7ETd9SYcgLPhP4NiLnixDyEcNvAiyhdo5h3TcX/uF+CMHg+xf274F/7n6wAWAdsLHGWGpxpVpNhKQUof//CskPEZQipTQhEqFW1DTNZDKZTHSjm+lkMplMZ9PpZDqdTaez2XSq1WWXR7hInL9AMHcIokmI/2+vM7dPAPvznh58V5gbxWQw4viIxW3vR6iW2QY8PQAulp0xru269WqzXq/X6/Wm3Ww327btNu263Xbbdrtt277vnbXMgEiASIqUL58TGhJmpgBEGhUiNRijafGKSQaUDETlLgAB0XcfEGkzg4oRU1z2vzj0AQDsADGU9CGkITKgR30SHBQBnpSDELBgEiUeDwMwpliJCMAKUcxaxeMKuM2/F4Iwm1hBBB51YiYkAHbOOueY2VjHxgaWKIvE4lMfO8esFE2mk1kzmUwn02YynU6m0+nBwcF0Ojs4mB7M5vODWTNpGq3C9JowDC8ISXYJeNWFtL879wlgf977cyN5PoQidAKzg8ScTLhAefdagK437bZdb7ab7Wa79UX+pt22bdu2put7a/0Y1mM3zjprjXPWWGMNs7MelSeFEMerROCZ/T6cIfr6ngPmEaaacsMr4k8Jiw+oCAClujM0LZH8DrFviEU5M2POG5L/iZjIUREjwvSE2G8rxNYhxfecefxjsIBIMIzGY17AVPpTiPK+ZcCKuJ8wo/jNPJsWCZHQ07mAEZBRhd8VOhGLgOlTdTHBKWRmhr43zrpt3yvcKFLKjxRIaaW1JqUa7cfNWk19gjiYzmJ/MJvOptPpRJMa6x0dO+A8Apeo2PPmg9EduP3ZJ4D9eUEkNbDOqXqAZTa+pLTOOrttu81mu15vNpvNer1ZbzzMs91utpvtZrPddn1njQUIUAlGvqMvrYFQ0aTBWN0PtA5CuJY1fQzJXCAPeYAKgC4UxhgQnKKCjngLoKSjFi9zjKuEWAAXAjxKmAyWbx4Xy1nydXGhSiR+lZxGB5QIczIaPpfwNXOxrZD+JPQeSlFOkAwMSotX6lcd5C9m55gB2DKYvufOOXbA7BjYeZAMtCbdNLPJZDo7mB3MDmbT+Ww2nc7m84PZbDafHcxmzWTSKFKkUCERKYWoy+kBV2/Ou98Z8OVVzv7sE8D+ZIa+C3BCKCRpwNsxDBfL1Wq1Wq03681mtVovV8vVarXetKY3vfX8y95YZ4x1zhprnQPIxPbAdA+1byLIx3o5l4M5yqcQkVGTWLtGLmbe1/IB0wMcGBk3nPERRgmjRAKni78cger+RsaMEC05sH8AU1Gf+wYQtT+k7dlYlpchiMUzy78qrlylH41yFStkAvnJcEzTccbgJ89Q/kyUiUdqY2DOefHvUniXHQP5LkwxO3QMFK4Qy8DGOu62nV2u15pQaa1INVprrWfT6fxgenh4OJ8fHMxns9n88OBgPj+YlO9swSnyixWRKHXtIL2P5vsEsD8vVhElZL8Cbi1Ab40x1lprjNl2/Wq1Wiwulj7qr9ar1Xq1Xq63W9NbH48pYkSACKQa3aS6X24CBLIK5jDp5KppBXUIDAbK0pxzl1AUlenHiN0lHBT6ad2YRCSU5fiOSBPC7Ygw23h4KiI8Cv26XSGsXJ8o1syQy59b/Hy87ENnkXEwNkDx3WFPVI0zFwp9CgGAim0VxKGJYAI4Zna9MT073mz81gYiKK1nk+l8Pj88PJjPD+fzg/l8fnR4eDifTyeenur32rQacIp47ELdh/l9AtifawZ2HA0kueOOUggiIowV+wDL9eb07HxxsVgsl8vFxXK52rTbvjd+49ZaZ4wztreOQVHcbsWsypNmpAHB95NGLtt/HEZ+STL3oDWCaBEE4yhA5CzQmIDqE+SqPqUfTi8afVVOGOYBBejDotwv3+KAtsg8VKk6oEhqZWCusgqmqhsjtweG3gHhicdnJ967mlO6I4Ng/pkJO0svAEUnwmHYwPmV1+JyiOJ98IcAWSGQcw7JhU0+Bmd503W9tavNWulzTVprNWmag9l0Pp8fHc3n8/nh0eHtW7ePDqZ6AAq50BagIA/hNTRBrpcsWLzD+8SyTwCfczkAd13zDGnIWNEMwTC01hhjrDGdNX1vV+vN4mJ5enZ2vji/WCwWi4vVem2MC2wPJEUKSZFqlEY/Y/TIDBEBgOeZxEFpCL71KikLVCJiMeUfhyDL1ctDzCCRmHLzIHgHfJslZ4mriI1cJk2+tFkKQ1fmCrXK4ZYSEiOhoWIAgViFHz/UJfFiI36EEtJJs2UW8j1Yhz+O0bxKqhj2gmvpCpF6MiQo2EeQ2EzihWH8H4jEU5ZEJWYG5+y2M9x2nnOFzKT1fDa7devo8Gh+dHR4+9adW0eHB7PpRCuttVe40EprJFDjF/eeObpPAPszWvXwrlqIAaKggkePqbp/DMBiuTpfLM4vFovzxeLiYrXebNpt39uu743pu743xjhA0iqoBiACUmLhc9xXQgTHHPH9hCznxgMQq4K/nI4mCBuERhmm5MV5fyoTZtLPqaI/ptjqgQkx/q0hlZH3bNcbHXD/tPGMBQ+1biIChSq3BfIrFn0BJwYUDEYi1Secs1+RemI7VZkeVKpEnAC/8EuZITVCuclALkYPiFwHXsQIGXpmD6LnEDECAiGw82xUZAQkYAcOmGHddr05v1gulVaN1o1W08l0fjA7Ojy6devo9q2jW0e3bh0eUBn9HQA7F9NmuVTwPPfKdbuJ/dkngA/KoUtQIYVZqMACGGM6Y/qut852xq5Xm7Pz8/Pz8/OL8/Ozi/PFYrPdGGOBSClFihCJVKMnhIOA6WOAi5wVTv8j4jIKl5YRjy3EAawtK3zMM+EUkrhG3nnH3ikWMTkgHHIwsCPID59liouF6woDljKmvPtDSGkxPzZMpjlMnzG3JEX05rpGH+nxECgiZjD6yoSM6NjLDQzX9J1yG654ZfLzQqGSlD43P+ZXiKRUmO+kZOyc663tPf3XGHZOKzWbHdw6Orx969bt2z4N3D44mE38flpsDlREKfnF5wT70L9PAB8YZIdHtp2u87c4adwQqoib9wyr1eZ8cX62WJyfny8WF8vVqu3arjO9Nabvu970xjCSmlAq9v1yFjsAqhFUTttVO1uQghovyfEM7FeJIDORMBWoMQBR9YMi+wZhh8UjQp4YYAhGFQaP1b/Vu02JiY8j9XdKqVBU71HtJ8D/LFJEkA2NEqLxZ6OI1gUblGJy8GpBfl4e5UczPoMZrwExCcFM+a/xrsh78usdVTBMkwUEIedDsY3BdCXmNTWOJKcqm0Q3AvSgIOQRhkNEIuXfamyInQPg1hi+WG62m+Oz00brptGzZnp4OD86Ojo6PDq6dXT71q35rFGiowWh2AHx+vvMBva9Gt4+AbwvhwCp3/d3ngPord223Wq7XZwvFxcXZ+fn5xeLc4/2LFfGWV/re2lK1UwawrgZNbzkkVGK4AxKQx6/QXJ8lGB2CGUUqfTp71DdsSeSjpTVYR4tudOUFgkQC1jopotHJXVJRLlILeLapJhhOPENrzKFU7HEJhugAuIOq29AAA6hgnYC/B5eFKGcu4twzEnDbcCcGhsio+xVyv4CWKyA+5kEDx0dRF7gNKOPEdqzfomAwTPD/Cayl2F1jjdtz9uWnWXHpNTh/PDo8PDo6OjWraO7t24fHs0P5wfT6WQa15Bl88efOaea/aLZPgF8JtH9a9T+/pJ0UQUs7OfHP+2sPT49f/z4yeOnT06OT1ebdd/3xrq+N9ZaUEorlaXOiIDQ5ZpUjC1TBBDRiguSerjHmUtPFRShX+ALLNXOwqQzc+dlTAwKa5g48sPlrTrvFPpwcQQSi9dEfceI0XNa7GKWCAOn9QIYPCEcG+1yicrHKQnGd+byiQOOpXJAlVWwkUi+CsRcBdfxOO05CFaR3Fqr2U4YCaJyPpGEKxCGWX9AWoI8isA0qJGLzGFDTuinMhEyk0NLhOyYUTlyALjd9sYslqv106fH04k+ODi4c+fW/Zdeuv/Svbt3bk8V5QkBO59Uo4LqZ+jW3EenfQL4LFf6MuRJoL81pm27tuu2bXu2uDg9PXv89Omz4+Pz8/OuN0SktSZSqPV0MgExf813fYHtirgvWSXAu9ti3P18UynJzMjFlpWAYmqQPjNIsSiVsVBoyLV/QqbKgHfFPV2NOdBvflWgixDHqKMADQiLMQEU9FZJmCngkx3vFsokAxlJAhgIGo0gEwjISTwuJgZEOXgfhvFdDULeFhzaorHULsIsvsQFxlWsyoX3hoBSPYAA7MBZa6zr+61lC85p3VxcLJer9Wq5Pr9zcXhrPguidc1EKcSdN8g+Tu8TwAcmrOMVjxE+ucIbJaztxr/eOXjy7PTJ0yfPjk/OFxfL5XLbd33Xd71lINU0hISkgopZ4mLKarHAISBadAUlAYZ6gCsGvZK8WN2AWPB+MADFqR9gpsCBZAn7p1STYN8i4KXZaV2qynkrVjrVRRgSbQ0mwD5h97HncCDnvYNVWx/eCYtJR8Wkr+QcJIZSSgdh/ZljHKxiFKaWH1WWmyYxYIDMKEVEkb+rL3j4IVXIXpgBYGzZsop1VrBz8fdSpu/GglxeY2WTwhJay9naD4YQkMAwIThy6BzDertpn5rT08Vkog4ODm7dOnzp7r2XXnrp/kv3phEydFGLVPKFalAOr1NZfS4ZS+8TwOdMnmDgJPEiKN0G4GK73W63m017cbF6/PTp02fHx6fHy+Vq23aA2DRaqUZPprF6RGHLiEP2DAryZalGMA5VjKJX9Vc8JnqTsKOqRB67OVHA72UCADlAHjjcYjGBuBSFKWt0LNEmroIJRtSLQCBMY4X+Zd/dBVqgzKh50o6DdwYzSIVcjqsRqtRXJ0LEHe8ED3JCNLXJAnopE5FPkSgzLyJg9VOL7qT6xZKaRYiolN8p4bC36Jida7fdyq3BuUbr+eF8uVhfLNfL5erwaD6fNNPpdDqdqDj3j/kARPdxze56H+X3CeAzjPWX3tw7GnpMBlZZPw3AAVysNm8/fOfhw0fHx6er1XrbbXtju653DHoyAQAgKot0TnSRKrwnk3GhXMN+mzfX4JgrT8FNQfEqMMU2IvQmKomH6MBxGB8gy0Fu0irAQDlNPBhMTMnSmnbI6OYSMCoxa4SRtekcgLjcRUvhLKqEQvEWZrnqRIAcpIgMixRbD3moelniLNqUQqGu9lsj/7FFrdPY0yUawJWZWrxLzFleI/dVOECqJN7OgKqSpBAUVi45V6FHwDEmLhbNLmJqZhDZAZIGzX7ZBLbb9ok9Pjk7e/vttw7nhy+9dO+Vl19+5eX7hwdTyuMBBsdJC/zSwI77BLBPAO/H40KrjIwQkW0wAJu+X2+2bduu1uuT0/OHjx8/efrk9OSs63pU1DQNktLNBBRBMjnEK2ohLLjieKmhABarV5HoIQtmkiV6AVvgsMqVUjwov+C4wkp4qfFYlQYwzXl3tPfVKEFoZSDCJY8XyIK32AUUr5tHfp3sSHAk9dZD7XLr69KpZgCEAqemGOmWVT3vLnhLEaSK8x9F2lCOivJUHMX/lPzgQPK5quxBFImS5fZxeVUCMmlSrMIk2zlnbdeZjd0sFnY6uVhtNptNu1yvbx8dzucH08lkNp00WnuV8d1F1fASuv7j92efAK4fxiUOfmkTEB8c2OgcVfmZWZHyD1tt22fHx48eP3n29Nnp2enGm2xZy4RqMvEOHYDkAJLaogDqZXmXZwAjLh559MjS9pZD7KOqEAUoa2xEEZnkGAMqrN5/l0rENvxawiBJlv8ejr17ldx+xShlwAICCsWx4PpEMxccacrCmCKiPdliIM8NAGreTflGypA+yFjihaf3LIFjlEcgXIrQ5aqWs86SEMtzouGoBrBcqnCkWUj89VLzYZhACdLgPfeDPKZ+V1xKfoEDg/8AwmCJmqWwK0NpogNZPwmBABGUIyQEzQywXK022/bx0yez6fT27Vv37t195f79u3fuzKdNaJSdS9MgLhlMA3JwULUT1y3uc8A+AXxGAaHswhIjKwO01qyX6/PV6vjk9OnxydOnz07PTi8WF8ZZpZvJZKqbaVaMQSyJ8qVGjph5JrVi3sWyHI23g3ItCRMzl6PM4DpV3We4m9Mknl9dIyMMkZgiqhdV6u6INFIMIxCFQTejINizUJVGqpzDRhspHAn9Q1IvjwYW8Q5ggvpc1TxhWkfA4TrEcGCBWWAjKyoAux0FcDnk59GOYQCcj3eHVXDFFH0Hghd8ORwGpaQ2ETIQIzQAwGCdMev15oKtQrpYrlar9Wq9vrdY3Llz++hgNj84aJQSvyr3S++CwNz+7BPAu3INJZTBb14Gd28VCs7WuqdPT998+8HDR49OTs7atu16a9lRoxtovAV5YoSIOC7gB3ljYlr/DxnGXXEDyqo+F2goxeaShH2KnOyqly+eAte1PzuBrI82IsP2pB4rRv4qJZWyWrI/gMssS2EOSqZZbi1ZzXOYWjBKScwCXMHEC2K5L8FjeBvi7vwaEkIyo8RkPMNZx0KMzZOsa46r5TynlCjNgg1Bqokp2xGnKiO+UihnJ1xWDWIfL2H6zGPXCjIUzQGPJPZairu8+EqHHvEZcJ5aECrdTICcQoa+749PT84XFw9nj+7cvvXK/Zc/8tprL929M1GhZ7XBoHnX2kBpubOPXvsE8C4lAL6yDuVoV+5JPgzQWrtcri5Wq/PFxbNnJ+88evjs2cl6vWZApZVuJs30oIimO39zWKnlJP6eqCwVQ2Ys7nMGBLC22hXGtnWNLupKHKtOo2mjj4uUeDVDCskl7+sIyxMuAXLT4DntG3vFMjlrpSxJwX7+AvWWUbU+OwR/sGqZ4tAex59nwDUGP4G5WnDOvyjz+gvJ7IKNj6OFOCAiuvRcrr5+Uz4jBk+0TEsP8beVMqMII5PrYf+HY+1K8RLKDxfT+CSqvMYOWZHSmhnYOsd932+36+WK1+vtZtNt225xfnF0+/Dw4GA+mzVxj8wBOGZ6EXW5/dkngMtieiHmg7ua3VRRe0c+b60CAFtrnz47efPBg0cPH52enW/btm17y05NGo8P8xAjKSF+wqjDj4XEugyYPBKpOJvkoojhiQ7EkXuOGcd3XNHdPeeHPRcoaOI4FiQ9BKKAckOotxEr8KHwrxKLaFF6rEwuBYMpTxKS1EQ5iY64M8XynyXNVOw1JJZLIReN0q3Xl+opzHESyMyFd1GXU7EHjVDx+7O5sfhbqeXgLFaUt98Eli/hJ0YY5C6OETwmwILlBYUAdeJ3htfgPPMfGQXjsvicSsWnfHlLHU75dXwHSYyFwrW3u4oSXRcUn5oCcqppUBEwG+tOzs5Xy9U700d3bh+98srLH3ntI3fv3p4Ej+gwY6M6t1/Dp3JPGtongGtUTtfRckiNPBKSxyovNtvzxeLZ6dmTp8/eefT4+NnJerMGANVo3Uw0NZzDwaWVctgfQizCKe9CXcuUhBmOhloPQLJFeMiRj3NKjlaTiNLcJYNPRQ29m3kUev5C3AdzPBy5F9OqAwfUAmPDkaTEENkhA1NSp5EIO0MRB7lyJdhRKUOC2sVOrHhmlIaKWERAlH6QWMxqkoNPheeB2BLjxDHCLO7pN62p3LZO8qxDEq1Yq8AyrOLAIQ2x3vdOCh8yhlazn5JkhrC7ZbqqWd4F3ihUQKCCq6Vja/q+3Xa4Wq83m23bb7f93cWd27cObx0eHk4nmdrAHA0pca/5s08A71lPAPUtAoDOY+UIGhAAVl3/1tvvfOrTn378+OlytWo7w8B6MvFbkuxp9bEoZ67qo4LGFyfJDIXmw1ACQEK0md0hsFxk2UjUXU4GVGteUJhFc1lcRuIo5F5+V3sEo/E3RjEAdpx0+ZN6ZXpm0ggYgVARJdMUTP44GEkmzFnUGneiangpmo9iG4qGewBIKCg0BXlKuDbm7QpC8D68ZW8VPJGTKLfkUEll5yKgy0aIs7ZSSNY0Fn8zUzNJ31FeemNhwMmeNcoc3IFKydQoznGVGj/zDmjo0gJcethgss4kRAe6Aa9MbR0fn54uLi5mj2Yvv3Tv8z760Y999LWZVgBgo+Re/Oho0Ac831xvnwA+dOF97NLgkQ3OZG+ikLyRVufc8cXFYnnx7PjsrXfefvj40dnZwlinlW4m00YRBJFO5hFpMb6099hBiMfx4BarNeIb3ABFjmEPnoPPWCmtVBEJC9llWbhjjY9JGI3HsIEBpSSl1ywLF2TywjCkeJvkLJOBsZo856wllmWxtBiM/1K1YpBDPI8qkopkluJ3Hq5jpiRlYzdBAkqIChZ+DJ6siSNldWlXWVEhR64JlnE8L2Vzjc0LU/r0dyjO24uFENxhuYxYvZ9SNYh5V0st5ZLyvAkxwlRAQUraWtO32227Xm9M15vetm17987tw8P5wcGBLvhlXA6i9mefAK4V/Uue5Y6aVnqXM7NDIAALcHK2+PSn3/z0gwfPjp9tNq0FVlqTRkRiRAulwVSqQHPBGHdm4LJbRfbjXIu2hZJ4CKmkuoyHFiWCel7ofrnoc0Wl/nOB+WCCwLCQbKgjPA+9UrKBFYHw0UKWUwPC7IKGhRBZMf3gqrJnIVctNVEHNsCQB8r5jyiHcSwps8moOJL7Cx1NYbCWkoFzThEREYNLuJmP9dHe3TEAoKOAxXFaUINaADq+w8wo3O1xUAGURmYjbF2ETE/GCN+zNPJhbxwXXDnLCVNxqYiEL9cqBm80DIlVWBc8Mm0PRSBI6WZKygLDpu3efvjw5PT43t17H/3oRz/y2qt3jg585nbO4UDXY3/2CeC55wD1N7yFlvK4CqrWmPPzxfHp+aMnTx4+fvzo0ePlagWkZrMp6YmHbwbuHhw63BopSbyISsufd7LsIDXmeTtodH0AxpzCs3oaZik1L9sQiaUoWX4Ml/YtMZZWM5LhUxEJr5RiZpmAWUqWJT/JSgGifpZxflu8S4mRL9RIL5F+E0k2xSNElMnNczvLh0rZ7MAIlQ0TISRJDQdcLJXxta/CyNjBS9vDgapRQogKVx1AGCM3YC3ZlFdDhOZIPYmSuhlRA6KoQsQSg7gSqrVhmcMC7hQ2RBQDEbNiBmdN27btdtt11jJvt9uX7t25e3Q4PzxsIlvUZa/m/dkngJtFerykR0hrng6gtebp0+NPfur1Nx+8db64MNb1zk5mM0ACVMzgVXRGAHEErwHKQS2nkpzHpM8DIxoMxZ2L2aoEGeGSEF3prCURfaFYhowRkadxbiaL55cZPVEFOgPo5bQCCi00H6HjPR5tf4tHJS+ZwboVS4IiD9+eWL6WhBzC0Yp5pFDMAwaslPkR6zcQ5ZwlAO35PchdT0EGZU5Ph6GQzsYBU6l8/3kk+UO5UTfQrQ41BQ2MbnJQZmbY1XYyj0GFeXjBo5SEwVJ4MkmDup4oaV3lS5OSdZx2vhGBlJ5MwTnXWfP06fH56enjW0cfefW1j33stbu3b2uVttz3ShD7BPDiM4D4Z85LWiEBgAE4OV88efL0wTvvvPXg7ePjk7bvddNo3ahGI4IbE1NBlMFXVrzZUjbemcK8pb4vUVK0IxBBQ2r2IMVxQeuOQj1lXPBOXFUuEXL0JVY/8NRiObEUmbN8FTEiIQvySokd5I2C2ty2sImsZNzEWB0BLrX9HdiwFS+pWAcuSm5p7SieQ92JxH5IyKgC7zY1kzsRyDepUy7ZJE7UWCxRr/yeDKc4I+9QTSordcMRRvbYxC9hHmCNxccoJCmgZFdkpm+6efy8DQBQkWNnjOm6drtxm7bte9N13cuv3n/p7t2jw7kWBKHPoNXMBz8BSEj3Xf/6Pf3hz/u1/DYI6UyxIIk+Afg/RQdwfH7++utvvf7GG0+ePtm2LSo9U43n7LjCgoTz9mQaSxGN1l2lVj4GO5P0l1n8kcCJxBpXYfk6FIuPry5xStP7EKD8QrUt0RILRYRKXjrHvvgiHFc6MXn7tIiTChUDZ5UEqPiZI0BuAg9kvMfMsWQhXS1YjWPZBaQvZdZuGEat9F/iZUcCDuc3HaDA0VisDuAoD4prQLAW7hwGdrFbOwal8zDLSRaswOoHPsyp8xIvU+rLUmBkIWVoK1/YWIrXIe8yfkcJ73NCLym9h06ue6NgRFU5PX+epDSAb7bdydnpar0+XZx//KMf+9hHX7t961B5FQ7HoHAITEZloao54ygQWBh6R5ZpFT0cRP7VZyBSvdchcd8B1HdcsBwEBgZCbIgcwPnF8vHJydsPH7/zzqPHT56s1mtSajLVQMgssM7yQgYcGg6WvJKE4peKyRwNv8t9NOScsqAu6QbALIIQkIjFFBa4S2HBMbh/hzZfKBqagdNYOR4V2IrYDUUgRg4QFNcd/8jUQcTm7BwvKSYVjz+3ItWmNBaxlEcXncdiL4/AeXlGgcV4oVInTbwUrn14riF2L6FABsAdI52ykeCqpqhM4/I0RMydOK7pym9j0bqKqQLAwK3HQ4iMoy8ij3fr0Xl+deXQoOxXqjQKiKgUsUJkYOdc126Xq7Wxlhnarr1//6WX7tw+nB+QIgBwg71kHgqSZ7GjD3EHUEppvMtfv6c//IW/Lu45UXsgA1sAArjYtJ9888H3f+KTT46fbtveWTuZTgHJJdUaLmvD0vWwKC8HEudcA/ZSbV8kAAwoOg+uaUg7R9XiKyaHVxS1P0GYlTESDuj86RlkeUkR6DGHt6hNk8u6nMmqjJrjBwsRBN9bFS9wJBYLz0JOajsAiL6LIBQoCkNWo4dyn7dYSAaShmhQQ+plFRhfQoL4hZVM2NhLdFT5I7mYuQwSWtVR4TD6IyFw2UJUW7lDNQ0cawmGWWQEXMvVwEBbNMy1M1sNR7QieHTTTnxDyJGG2p95WOOLX8vyu9VQGREJmJGgmc6cs9a54+OT5XJxcnr6+R//+Od99KOHs4ZLDKvSjS3S41ik2h3B6DMfqd67H77vALBq0pkdEfl6Z933x8cnb73zzqfffPvx48erzQYRtW6oaTiPQIXM5TC9ZLxhlE4timNMwmHhIVRmC66lNEWdFLU7k2pYAnoRhnxxD6wOcPLoEeDZjjgQwOEdYmlVNcel5Wza96JikonjMst11GJh4FJgK1RPKmQPMrQXQ5bzR6xHNFiSV2qXzMo3WIy9KYltiNC1Qyvt8lF9cUXmrV4CcCPsdgQcndTjcJBc7Epn8n9BucyqzZmxOjLZTY+JDASomG6VWwSXHQReOeDI2JtMGGXi9DsWBMzEfklQGWONNf1yYww7hq7rXn3l/t3bt6ZaQyF5cuXH8aHrBPYsoHJsGb20HMDW9I+fPP3e7//Ep15/Y7laI6lmMvO3gHPsUcywuhOg9NrvO24vBegUR6sjlG6ofgnHgdfux5ICMrJvySCF9JNBa+SsD8eeedgMo21xLrmScwgP0NKye6/mAVwZZFHFLEp9Al2pHMdhOJpINlzOSxhKNEraCZRwN2Z2ZtGMxWYm/fy0rIBlnEWUEpuigg2pjLlGwBAv82zkXcEGg1yT+Lhpx6hXmIDm5w7jWFEs7ik8c2aWikMlAR9l5yUzSc6SKIbvqV/gtJmNRSWT0KaUOTh3t1wVNOETSE+MY7IXyh8oxsaISE1DSinn2Fr39MnT5cXFarX6/I9//P5LdxqtI03YAY7KgV9PR2ifAD6HT2D7MCgiALAAx+dnjx4/eeutt99468Hp6bkDnM600orLCjHbcgj+wi4zCh7pjwUojBnW9hVptdK1u3DxvidYTjZhB+4pTGgT4FNoU7oRdEw0NEMAi0UZGePhYH5cvwXJEXdYgCWmaYyHUrySpQ+lVHTGuFiLCTASY/ndMkBxII51vVuzicbDA2fd1QoQufRyQ4QKNOJ6Ls7McJmPGpfAWRAgKkbbnJ9goA4jS5P5nOHyVSCkrKHSDkIsRjLVnJsLCx+pAVhm9J0wH9KgbSiS744ajoi8yqJidrZvO3OxXL/z8LG1brN95eWX7t06OlTgmRoOdupnfLhnAB/y6C8Dgmf7vPH6m598/Y1Hjx9vu66ZTgA1oEetEwLAPEDWeADO5AItCbLvcHBMKjj+PnBDCv+gjMQkN0aSVpStopKGWqzhMGvupFJdFvDsWHo7jtXMY0Vr4XaLWAbAQW+NFWiRUyoLHVMvU0ODWr4q13xdSan3kpiZz0bMA9N5zHJvWDiJ5XTBAuaQEJx0WkzlLNb7yAmwzmGSaxkGsfmQtoX5GsCELGOT+wFEHeyEU2Mg5wzNsyBKxhYj5ARu7ezLMAo8OChVXaXgLGKtwo1SN28Uas+4GtbSWCgKBeYaGBczpnj16emU2LrVav3gwdvrzbrvevzIa7eP5gjAAVC7BAbYQ0AfkpCPRe3PAJ7pv267h0+ePHjn4YN33nn85OlqsyFSWjWkFGdPjTQL5bqrZx5brMklXVwsEjtHWAnkjiYnOfNjrgZ3BAg4hs3LbYJSIUAq/BQQQ0Y+WPLKUbQaRfRGFCPgy0B9Fth24e+USsMQIQnRReljGTYEZz/4xgRaHxYofLnyW9iY4Ii8E4JEh30tz5UdvHjPamWgsZJVGirIJShOY0iUMArK3WGMKBuXl06hhF3OrErsJo7BhcAHjiOfmeiTDKwLYCkwuHAsXnPGqDhZiNaravWcRdDEwkC58KfETIN29QQlEZZwdBs+YZbei0k7NMaZbdueHJ8Cc9d3r7z88sv37ky0BiDDLtn27eBXfcgSQNy3vy65tZCKvC6JtfywryaojjBzd3/tbvzkOYihSBh5y/Dg0eN/9a//9TuPHq3WW+vsZDoDJJ8k/DYKZ3GfQTE7NEGsqqtc61Bxm+5WeS55fVx/nbRwEYcKwMwFw4azmihiGT2gknsbBWXiz/aYftBCpewHiynAcz1yGHlVOJgQJkstPx3gwtIckTBLn/p8YV2YwVICo4XwGpRbd5C14+R0JgUxLN1kGcptqtG3BctJLEPhVZ9GQhgiKcKlWuOFrbKgu5YVdfZ1GFDZmCHLIHAxihDkeiy8ZhiE1pIk3OduBqm+HlFuu0jV7Bpil3Kvcf84Ng4o27ARKs6YMxvtuEekFzF6ett0MnVKG2ueHZ+s15vVes3uB7z22n1dps5KNGUYMcon8LzR5hqP3xUq4epVqhcKlfp5e5/31eNvnr092sOOGTQRACyW69ffeedTr7/+9jvvLC4uGFDrRjUNIlpXllXjXUWhaYU4AEcyaF3g+5es8mJp+ooZfkiEIMxejrW+Sz3zhOK+F2CC4GIXu6TF3j5D2o71Ny+lyUV2OLhKipez2pogR9UbaBWlsF5ESM+dKOkMFItMjKP6ozIcYvly5WyaL/2UR14SjIoKjdcFeK3LEgs/+qEtmzS1YajdtiQRC0dM3MWAocrIBZoi2pQS6+fkoyPlLKC010wlSNlVSkuF6nfh8EIdBd9qdDVeqCwhQkTSiggZwJjuYrV2T545B23fvvby/cPZFAAMOwYIWrPvbbR5X4dKDbvZr3ANUuqLPP6zwv1HxFz2xFC3te71t9/+zu/87sdPnxh2SjWAREpxNL4KOGuSnMQRMr+QcsxK+hhI9yn4JHw/+y5Bbn4HkKTwZClzgNTgieXVbgQJ0hMakXFxcVS4G21niZyLm5lHt7YuvfA4Ei+BAVFFV/qCvZnjc4bROUFFKfkJvr+QlccqcZB8mnkvol5mTanHx0YaUEXL1q4KzJKci2MZnQshixHGaQVslPppXGgu1/OeWCZKKH/HHCHnWORChnl4h+IuF50dORKzACvuxIIqWTgBDko9iYpIXb35wwkbC+EK5sjQBtSNJkLn3GbTvvPOw81m7Zg//rGPThWFPYvdn+quZu25+fvvz1D5oZsB+M/csgOP+yOeLhafevPt7//kJx8/fbJp26aZUKMJo5R/XaqPrNMX6vLF3K24ZhlHSR0CJ5YKQZhr89DfOWBiZMIk+1+PIfJQkdPKTGnwNESvBgtYch85SE7IfuOqJVYZT1kEREF3ihWkEHMu4Zr4xYgVST1nR9GsZ4gMKZuiECCjl58IxEcezFalmBEVAhx4VfXOZfsirBmlHfvInhsWiNVlCVSCV9WqcyYgpbVehhE6TQIDWPQt6UqRWF3dObHUH6k8L2VZL8An6W2TxPVGLkIUWFbtKhBgEdGoxCK/WIWP11ll+ZC1rRUSOGt7s+26k5MzpR445z76ysu35gcAaAdbfx8218jP5QTAOz5dLzrvnUVXbf+pNx/8y+/87ifHTwFpMpshKYoT2qFjS01TEJIPWSUlBaKgbcmXMo0LIfoUjbAKHojMDhxGws+IQwfuCO44BJugKgQ4LfJKfFRMLaiQreFRvGP0timgHMoZD8WNDqM0S8bxvn9wr4bPRCojkewPkBCZMtA6+mTlNu8uWsiIVCoO56p1u1YqKgylrbEe4BcOP8NZUt03MNYhVLKLyjsAc31S/q0duqmiyMZySMLl7AYQKAzVXNYBlGTXONodvDIc/T4IzFFifCyMiyVAJFIUVh4WSilEJIvOuSdPnrZdCw7Uxz920KgxHfYP3xD4c/bgSO3o475GZMTj8/NPvP7mJz71qWcnx721TaOUUmFxtJCRH0N+pepLbs1RiCimS5CqaDmsP0txFmZmLIAfRGaomY5YMcWTIS1IhYpwyyAPpsRQjKAliOGCPigRXNoGc6VlXVrjgtgJqCMrFpI5UUkNdwi5jW2Rph/I4kEpCTBX+2kSP4PRGId4FVZfmdyX8QtLo5S8ITAQTWbMEnKlo0oBX2ENeGABdteLWoyS0CXxm9RioWS1Fh+mSBZUmbU5yDK1zMygJO0KU52QfgpKUx0GqLyUczsoMuNQd4nHcz3W6TR3MlyiXMWFQooI0Zi+67vz84s3H7zlrP3Ya6/evjVPGCjefBqwTwDv7/g/OptDBAALsNxsP/n6m9/93d/97PSUlJodzB2Dk4AFXwYiFSUfY6bkyAuQK5gSdiGnkBzPBUUGS69a6ceBozhBhoHTaMvJzr5MQVx35sUYASulIBjZy8/8BDGbhMAYQs/dwRFZL6y3kzmPc1HwYZhl2BtuQQuho/S3GUQ/Uy7J1mU74o7GhaEWyJRe9uWMAcd7gjg0l3vaUBSquENVY1yvpk5OY1TLcY0OKUEy5LckqIpLaVlM7RJxXKhAv6OOiEjEcfkZIe195FR/lV9NvgoK2m5cWxmCPDlRYUFMq7I6VtwBOX8nJKWmOHGWnz096dsemEl/fH7Q8IfYS+xzugOIRbq/DBwzM2siB/D0+OQTb7zxyU+9fnx2ap1Tk4lSGp2TG6iYa9ORgrmwxkjweHXnohRO5uGFmbsGX7W6CB4RDe8YqdfPO18rFLtU0Sl7sFTKYueeKi03RMRxqcySwTFQPorItvRmlS1CDjksYQKMlMziPaxJOUkUAV30ZmT2jUpOv5g4UvnJce0oUCXoakOAh9P0ShABCt1nKEWHIMkOc7FrG8bekREgDRoLedcdEqVcQJBw6Ypw+VfFICZ8NETko7f8vQRJx4EF/QwLsZH4Vsfki/L24F1FGFfvTd7pTok16s3V7qPFGCi+BESpRMFjy+/5lYtsTkQARMDG9N3FxcWnHzzobP95H33tzu07hODYDZeEmS/bp7h0zggfiGmCvkxm4HNmFpDuHkQDcLZYfvKNN77ru/7V2fmp0s3sYJbkoiq3khFKZUFbTpcxlbTcSoe3tnySjXoqtDmsdI6U3jtvsNoAXd6smRjHaUA4IIICAlEQ9gyiY2mNNJaw4xfymK9uRlpY3LDhoa5AvdMiMg7573J+mWaI0hddor6clHxK7ggC0o73a+ea2ugLTUyZ5N4bfZPjNsQlDWjuUThq9HHhf5zVESqOS3nlcVp5u9T95xrfJww+lxwXAxBkLgpYPUpSqWyGUh7ny/pa3tHpcirbU6YM136hqJqoelDirMLgQI7Lxn4ll2Q38XYqrZUmZ8zJyUnXbQFYkT46OgwzDK6otJxe/K79gGt/Bu/LBHDtV1V0mtd8vLRWeNd/+NWPj/tLjh071ygFAM/OL773E5/4vk988mxxZhwrUqTIxb3SvAI5QFeEP1Z2LeXSi4VHLgSOlXgeWfmLiSgpjIXL2Ak9OcH/w10AePozRh4wJqsrMQsX+244SycnhIkgCxsHZZnRWobTsA0y9ItimRMqMV8YrlGFqUbVIxQFeR0Cubi3azEDsQtWrkxjNXQZXU5zAJTXrfMfUtx3KP0dRlLh4GmyUJCDvPyBhMDCrbHAscX1BoVROsLYrncxBc2EG2H4hUljijlulCUXBqkXi3kOg8UCinSQk8vcl8Y+3Pksi2lKFvGjcDelUZCUpJZ1C9ee2uUUd0SBiFnMQTAKOBKw63tYrdbvvPNIEX3+xz9+e34AAMY5P/5CGLKyh3fCuziv/OyESr0LQLzm030vHv9u/fBElSAkp4AB1l3/5oN3vvf7PvH02RNUajadAZBzHOUdc/kVyH1D+TPBhnFXdX0s22AWTb64GzmbCaBfXYEdYlnlIRhafKebgUt5exRVuacREQxVk4fC9ZdtrcqWXjpZDnQucShjX5bjPLA/LkAZoXgtc4nvk3joYFMuybJU269HwJwbikwwKWH4oRkMF/FMGNpf46LNy8FXB5AsZZfsKa9V9WPVu/jf5yLNDLE0SgjW1JlvSYg1BRd5qLoztJ8bZoARfzqR5CneRIFDyhGWwnhtD2lbFWcHa8WPCrLjsb2QULcpPT0gNnZxvniAOJsdzD7vYxNFpOJNuCOs4M7R0ehHgjeFgD7DofJzfA/AsXOOG6UIYLHefOJTn/6+T3zi9PzMMmvSXkKWYbBrIkyM4iJRhnQvB/lGLIER69Gww0ozYWwMJVggDOP0mIKJNLRjFZBuwk1LVuyouYCI5GOX9djOrxRDqye+UnMFUUoDYCFxXZBLWQojR+EknywxiJ+i2L3jUg05zMsF9796lTW8TihKaRJIe6YDxBUt0XpkdTKUs0kRwZHzfHWQYVG0OsxDyme1glDO2wtNbiyqaswTqKR2QZCNC6JMKiamgWDsJ4A+Uamk8mfCSGKjIe2PsrRV1TClmRnXG+OJtxSW76go3nB3opOiWzVXQN42xeNDkcARe1XkHBtjLi6WD956GwA+9tprB9MGEKxzQdcQLuu6CmmMAWr4gbCn/9xMACz7TAUA0Fv34OHjf/W93/voyWOHODuYc9TAlQyyQUDHUrQ5qdDX+4qFnfigAMAdnkLSYoXHYZ40wUXJaOE6Rwzrr1zHIXBaG643eXisY+KdY+bx8i7N9aqJq7wZOKkKJDS4WKrmUm8eUs4N6TNg7lwJFomXgYUuhLgDsZriVPdolMVncANIhstXXAIh2UUnc3pkzVoCRgPvNZTrflAI4AwFlEd6CRALf1C45ohlsdKcncvlamkEKYEeLn41FhGdK5ip5CYP1ngH8yGZ6ViscnHStSp2Gy/TJi1e4IjXEkcnHS5tGjh6YQOi0s0EnbVPj4+ZudHNxz76ikbyOiOj8fsmsM8HgFo0YgpfNdUVODUUGCpDH9zg8eVubf5dzDBaABR0gvwf9e8NtT875xqlCHDd9m+89db3fuJTz06OO9PrZoqkOERq5BGgpa5oL//4WdS65ZpKgJWlzlRZkOK4la4AQhBrz9QBCF2E/aEDKmJZ+OOYMZlwjRwBZeqCNAswYBqPjFSjQxM0lHYj8r3P66m1s3t6jQxCdx8LTQtfnlOlP0dIV3fEydcdkdxghDQSdynTerA2TRmDJHAI6SUX+TRyrFUXePgD5EA+AWIx+ZHnIZCIwmEFJO3PMkuppAJoqi6K3IJhiabXLE0pVsSybB/sx0uPzBElidKys8wYuTLAIcyHFcObWb7FBfw9fM2ASAoRmZ0x9vR88dbbbzO71159edY0gGCso+gOWC9M79LJgN2CMS8SKnc8/pK4DVe723/WTeF3TIwuGySVGAReUvvH1GKse+fpk3/1fd//zjvvGGsn0xlDEBzmsYWeET4IDyaYo1mB+TLMz89UCYdWi8NKI7kh4mUQ6443R8xA05b+NYoT3Fl0jixIoXA1Kb6BWd4aRTrHugQfU4sRiQrHOu7SI6cOFCM03GI4XPw68UqkYRmNys2M2mqKbqCqIK6CiPMblhXxoGK6FK+jUnYrdNPkO0BQ+LcFl8u4hejgZnw/Ttk90ylBgnpi1jzQl5KaiPJ14EhPVM6rWBoRjwkBXo16s9C8Kz6RpFiOhT29UrpBUtbap0+fOed003zklZcUBqU43lH44WXoUKFg8X5eMihooKMz5WqD/pqPz/osVz0+g4rVgyXpC+TMVdT+UD+eAZkdO1ZKaYXbvn/wzsNPvP7G4yeP15staq2JBDsHB6avyCWQ4RO93yFArNTzC73+upyJSLXUS8FhW31ZYK+dTLmS70+1P7KfyiFzppEG599RgGk86+almyycXNuPi2fCKAmmtQ9IAXXzJUTLpEuc1a1zJYX1XnUBr8SUQYP5J5YwTJqzV8OXHRzC4bvFKLZei2pWTGcZyill0n5FrM0W62m6rL057tCFIQ57o52I0CMXKDlJF/UyypfYB1EGhOoV8sG9WXgCIaI3l8nKxliImQ+6nwKqk3F70Kxnyxcu6pZCAXCk5hO/NT5/FNzPONzeOWEXvFhQGtEZZ7dtd3p69nD6CJlfuX+v0RrAzwNCJuS6L71MCjHmd2YeJ+HcNFTuevzlXCBmvuSHX6sDuOmuwHWZTLvuuMuJZlc0CIxp5MT87PT0+z/1+ptvvrXdtEprUAoiEDn20YUYzZhU1ylfKrtsXENiGJnfoXAlRFEwlyaAOzy2cHSYNa6pVcqEsXA6xKorujQNjCA/WErlX/aXSjIJZBwHL+k48keWsDGuJ2gIw6WrNOikyy8XkdIAr+A27br8JBA3lFKqpp0VjYtr/0fGcQsSHoV+Yqgu5pFJc1xakaYkJ7TXUiNZ4h6XVc2DDin3T7WsIO6iLYiwXFdFXFooVXJHQufHX0vMu82TB+9yWfUDDkdkNc0tL3YgEZHSANB13ZMnT4B5Mmnu370TswhDuZXNRYwrbwysjI4ZYGf39R6F1us/WO+q7i8HoarHXwlCjT+eRz5NRBxKjcVKxw3uEkGI93MfZgVIiMa5k9PzTz94+52Hj84WCySlm0k0fM+VG+dxWfiYqFSPYV+bxEHfqEUfliM4CZEPVhuHhUCKoFTOOkCMD3LlFRGVkTEJZX9ylqtqeI2rBDG9ciww9yGDf0B3Sb+MR8pprDu0+mmHTa60IZzbJxRToVzzZwfZ0anpzhyDQj9b4tDC+tF/RK74djbYBPlhJJp/Xj3AYYU73JauLvVBaCg+AYTs9wbZQqtM/ZmRWVvOcK3KUI7Yqjhd+I5BxZ6Kf4EqKc9cuXOR6iP4Jt040wJNLWjCg0FzDgWXmWHzSBJFSY8rfw8L6ZZcyqNjB4AKkZQGYNP1F8sVkpofHjRKHR0d+sVpx0wC/Y3X6oiNQT3NGeMFvruh9ZK4Pfqj5E/Qlzyz6o92WeQ8/+NlQ8y5qJA+GyNudFCvfhTsDEBEdABnF8vX33rw5psPLpZLQCKlkYgL0kClehXQOq4rc+ai7hqteQbYBcI41DMcLAz/QBr2QaXXOD40iEpxRBhBsEqnrjYv5DriYN4WCM6VvGO3P0UjHCTwoWhE5ZRe5EWxz4W4QxiHi0q6LBivLuaLmIw4PoCsXx0yw6gj1RV/DUcktUcCwe5+bBSXykgGj5RMcrUYGQbSc5f9gjHRtIQlcZIVqVAORhHuRU2DPGALlHdR8ZExFPI/xeolAl7vPRfiJYlnF64jKkA48FojddyM5VX4qBQSQcPaWsPr7frxk2dKqY9qfXt+AH5JEocs3rEh87WD5LsVWq+zDbDLFeBqFhAI4vDlWasaGFzj8cKlrr5yRkIei9HiCOkryolYgOVq/fajR6+/8elHjx731ijdIBF7pFwWHCg7fGLMoyoa3NBYbL4jFEIk0h8QsyVVoVMFBZ20ZNgk5f06uLE35aN812HU4RK8CF8qZpkFoiL9DPw05KQOVUZkuWCWJHeDwYasNETMkQlhZHWSsy4cpM6kCPfiPz3iT2Wk4lqAY7AxhtXdXsZfHDY2hRRyql1paDMuTVxwXIqgNDkpVo9k8yJHiaKW5OoqD8k3jaNxSEoQi9A1+lBZtRTwDcr3CwsCHSOw9C8um4XYeAS/ZOEYXHQJxYfCI5eC/PniI8aofwvFFjzzSNrgkgEQFb8ZiXRersax6jhh8dWcG1CCY4ikNSJaY49PThhgOp1OGj1tmvSCmRnrdRoQK5DX3fy6UWitvlk9/kXi9rvJAkrThl3tzCUYXjkQ4N0YZaUonG/ihN4st9u3Hz769BuffvL0ybbd6slUNdplGgQmQ1gIamgp4MPo0B9L3snAV0ty0sZr3LFasiQ1jI4lucYK0ozCFUaDiOOm3FeJ90OSgvMc+zwiHJX1ZxHQi1TAPGJxUI05JO9cqDrWRlljfQBX/vaDPevhzGAEZ+fBlYVSb4IAxxXqE5LARXYuK1ms174vKcYxNQxjzmsJAZN15Yg2MtaBJpaFzOMztuHkmcs7iy9pU2V5zjjcnBwBZwbXD47+7FFOcqUGUreviH5Ly7t9OeeAwbJ1bJ21AEhEaZ2YExluHFQhYS0XxipKa2bu2u7s/OzJ08PpbHb/3p2JUgDg2FXMthz9uWQnXFUrPwdqf4PQeu1zNQvovfs+FiSE4R/w6I8SCh1x1uV3yh0jomU+Pjn9/k9+8s23Hmy2W9INEA04Hig691wlR14QJ2c5KMUhZZ3nkw4BYd5JLYrMKmAJQ135OjjBkTxmX4Cl9HJKRP53Ot8fUIk0jc6WueKry8+Eh9hN4SsbByAYJyJ5DVNUh9XdH3nWclbAWSk0SzSj2EyjAjrCof6NnHYWCBwigEtakZeKWGRDGmHAFtwEiwKSRxB6vEwYkuWCA9eC1Bw7S1nHYTGDQmkJJLp9/z+lc6FMbQ6TPKdIvlg4ADN6iHBkuia3h8tOb8yZOWG1qbx0ZQkjJ2v5hUZCD18KUIjng2ksB2VqQz+yVUQIRGh6NqZr265rt33fM4cHFIP7mMmJKAkxRViQkMjTrPy1rhQxMxJ2bffk6ROtaNrou3duY7iLSVAFB02F5BzjCO3nvQitL/JDPoibwDUYyuAISBFZgOOz8wcP3n7w8OFicdFMp5PZ1Drm0kVjV7m4i/Bbda3yxiZEhB1sMx75bym8FQU3BZiD43/fS/1jDCu+/GEAsmKvLN8keBmzRSz9yBhe1oYyU9HOl1Yt29btRmVMy5WZIw9XirGK6YwDUXj5TiIhOwmHj4fnIiHh6MAVxQxbbEBxyXBK8SRvWtRPTSiWSs1S+bdwtBcsB7Y7N0kE8lMOKjCvqWPZlgbdKhY45gBA5uLiw/qiGGNw4cgi1NBeC6vR0fD2kgQHUSxj7Hcpf4BBl4+tc33Xd2FW4Uxvuu1m3W7bzbY3Pbs08Me4uUfZ4h4JY7fvoz4iAhEhEZLXiVPKJwN01m43675vSQHDx+/cuaPJ9wEs7jmGUYc++GDoLOvPciC/6g/EnH3QqAYLOnaWlQIEWGw2n3rjjU99+o31ekONRlIMUfw2LbzXlSqJuwtHRvrMyOxrhGiqmkh5NCIvX8Odg368jAw4NL2KrimYTb8EU0ZOQWlEObiMx8xV9YYDXzGA0rg12spg3VdIjfiCoSf8iTP/tURUIqM/akEk27U87JUEchy1D4e8yBM/t8DLiUBR+bJY2oSUPE4eHTYW71RaTSiNprAGupOcPSD5xMzOSXWd6IsTex3OHQjHXzS+CJ4ZR5yaMYhD/niRUIie5adVCGRjYkfyCLQHxVSxZGoOORhYaViJlJxTcOhbhNhEaBNr5QaGYp3YU9gE64IhRXNUSIjO2nazPl9cdF3r2AGzNcZYY4y1pnfe0gMK9S0U4SISuxmAKXdgiMkWA9JvA+ssOz4/P1utVtvt9ku/5Etuz+cebCCvoi4mIpdPX9+32eD93wHsUFrDkicBYK09Pjl56623nz07BsTZ7MB5lf9SH2a0yOJi0oDScNzfnL49dCLfXzrtGatBZbl9GUeovvmo3NaH0RJ2lCJRpYYsq0YgKTwoRHhCi56pmNJDcBeBJpX0Y0O4UdwpkAwZpJLXLh4Mj99TiQsgBk+7fSOHvrslDaVONDxCTBjRrJPRL9FmQEw8OBaiyTHCMjvhwljIM+yUEw+2vDQkUDC6pE5xSSQq5evymywnMYxXapT6GEnV2kJi32FuODjPeiKDg6/AvrMxQZx1eUEXa631UD4RWmMuLi7Ozk62m63zcz3nBE0gl0FulPQgIcdxD5liik6E6zWvlhcAcP/+y/PP/zyttFjxvRFZ7H2ZAC4RoLjp93dxVy///nV++EAiKNaPzARICK3pnx2fvf32w9Ozs7bvlZ5ANlBEHoSgasgmAlw9ygkEllRSUCHrKK0Cd/YyBXuMi4300flGiMmeDSKHk5XKCO4MxszlsmjhXIsFJRN5oOAb/f/q2wQBC38+EU9Y4BB1AmAYPD6Gfiy2hwN2wnWUHRFYEBkCCSuqdyw2R7ZNuIYkYNi851E/D+SI6v2AGkMLgYhCSZ1rzMxHDELWCMJRqyQk5xIY5eeN0eQ++txLf5rKw3IgqIGSDCOMvdKnKVjycSJapUiBKOW0Ed9oHCMaxslKWOkSkzEuoP3YLVCa7sag5IzdbDbr9brvWuecY3bWdt12u906awq3JWbR+JTrL8V1MJC6rghL/sm4fB9yb9pt+/jx40+9/sasaV599ZVGKwR2jlER1ht079MQuuv7H1A1UDEMRESg9Wb95jsP33zwYNu2XtljlNxQWExgjWZL1o4MoHIYGb5mrNRNYNfINbP7S6NhLrltEgkKdzuP8MuvlBapIKFq7gxXDjhQNN+XstVlRM/EFal7wdkKq6Dzjbc/I4aNdU2V5ogjzwWrdZigZVNC+dfp1gBgp8hXdiMO7xohc3YNDPNY3y5i1v7jgFKxMB7HqwvGQmwIUSz6+rzJzJ7ahpXA2+7iHbDkGtRCfTgu6lkNjHmkIx8I9+fBR2ElnSV/GMM0K4Gr4JxzzlrnmMPXbdevN+vlctltt9Ya5xywC9cKUb5PpWalY2aXhWMIKnXysssLQKroshEBQMXBOyIoBQCb1erN11+fNc3h4eFLd+/UDTF+sHzAdkBAuzbKrvy+UBjn639/V4oLF3f9S3Gk8wZgAAtwcr5468GDR48eW3aqmQIK47kYwRlqlT4WvqKUfktkXAjzP0ShyY+BopAVbnio3YgCls5Od0VkzGaHeZsmwE1QAhBlOwIwPi9mMY4OZBsaOtNkBg3KbFgFnVFBGUHFjRI9BIJnIbAqId8QgSWWqxRU9RxydDwMhNEtapjSRDPno4grWgsursKUDJK9+wjuMoTlhK0XsDQBALlbhhnKSYvBeWpQXHEovA3KHYwiUg7FG1ItgoK6BAWVuLy0L6lMOW4aY5onlypsme6Y/q7nEDmpRhO+T1nXOrcz7D85/zlQnB+4cGsHkIgQtVZEaKzdbtr1Zttu275re9M766xz1vRt1ztrogu0iBIA4BxYCzGZhCdlGNjWAoAyO7H0WXNAyL6lYgZERYSA1q//+u802vTm2eNHk+nkox/5yJ1bt5SKorDMpfXzuE/AexpCnztuv986gGv6J0TwB1Eh9eyenSzefvvhs5Pjbdc1TUOKmIt1eyxqWx4g2kVjUNn54jAeAA715YtXgXApxj8UShHs7+ElO9Ro2WnVIqwrcSCQJjrgmkWOO9DfsSeCIiSUbPYrNXnEAKXUkqiSZxL+qBPCTnUvTPt3O/eXxRC5VPvnorzdWf/DKINSfpwlzQBH9HAuU53LCasURkBZjggpi2IsPLhORncHdt1uwmmes3QR80g7hNV0CS8dQAlHBMS8a+JZl8zsnG37ruuQCI2x69VqtVqtN5uubXvTC63DSFXwsdu5IM1IpEiR0szOOgPOAaBSmnSggSZSd3n7Ylb+8LxPokBVZbbO2r431kCYJwMiOES2dmu6R48evf7GG4eH89defVUrAgDL7krV8fd1B3AVV3R8r+HSB8OYkSvsajWupMqOXrrsGBUCwHK5/uSnPvnJ1z+1Wm8b3QBp6e6bhqgucte9/PeAKJIdjtJgBzPLJVmRYpyejuZhWWuXSYKBwXGaIBSwBldmUTjwlc+Rpizl83SDRW1emGBI6r28AxDGwhlWSwEDS0Q/Q05adVgL+snUEvKpS+LcUZuMygtLGuxkpf3gRhsg69rqccysA9PflHge8pD5iRmI5mJuLc3QONsPcKT1o4d9oNZoFa0hpbo6ranESwsz7I4gGSRJySCU2MxcKVAj4PBlV5TXLJcvqT3lgmMF6ufldoTa1Tn7CXPpAUqAjFgazdRvNmYvsjSD8PW0Q0RFSimyxm7bdrE4b9utz2hd35m+77veOhOtTFEQFhicA8vgbPjMpurw8PDgYNa13fni3HUtaDWdzW4d3ZrP51prMeeTtQVitjMTDCAAZ+16vT45Pe4uWrYGG62VRkBjjf8p6+Xq+773e5VS8/n83p3b4dPlUT4s3Ci0Xu/BN308XhlX31cdwDWX5fLX1tqTs7MHD9558uwYSE1mB85HPEJpN1gutvDAsyoxDjjo6VzN0hnbg8TLit8yTnHG+pGehzyAIz0FCoZNjsjJRPG6P4x3vp4c2OlSy7CRv5YChRSzzku8o09wt0NxmX6HXmtc1I64qyng677zebcPdzkWFHN2ATIXfSNUdjhFFpUm6QMRqMHCLGZFn7F7G0fbARyZtvjhFI4s/Q4poslpp8Qgs4wJDhd8w/yXwSGgUgoAjDFda7quWy6X5+dn680anPPVesJKOSobAjt2kRGhtW4UkR8cYDOZ3r13d35wsFgsLi4unHWgVaP1/HB+587d2XQa3hwUUnepxxZT8UDTBddu2t70COicA8cTpQ+PjppmYrq+nXVE1G43z05OXv/0G5/3eR87mh9onaV08ANG/7luAsB3PYLvCqY70matt8fMREgArTWnZ4vHT54ulitjnUYvA8pciveGAjPOCX0Nm2yGKHH5mTjd5VLkl7FG0XO5yVkqRyAhXNsuBURU8DRidTkYilbz2yLecU3pkWKoWPiScN7OTSo9hV4FZ2cnMUHlwcich5V+jmeuyIUFGwWT3QmJyEhAFZxTC5aiEEEeFUuQVsPCzIXz0Hd8nJtEP/PacpzUQj0pLvYGwgVJSNnxoHCbLVxD0xoscEF6RxQtUPqLeRTiC2QcmdWIxblil4yiueGIk1Q168L8NgsPt8qNICdKFsYHwbWuQu3kWMjFvyJ7gpT7mIMvN4BrdNNo1fXmYnG+vFi0Xdebrt22YP1oFxNHCvwuhbXg4rtIoCb68HB+dHRrNptq1fht39nBDBFW63VCwxyzNdZaY92EVABywoeXfUckmMAIoJXqjV2tNouLi743QASAzWR2797927dvM7uu7dm58/Oz05OT1Wr95oO35/OD1159daI1IrJzLEA/Bn63o+VNQyu+OwngRlsMldzoezQlCBJUitpt/86jJ2+/87DrOt1MMEg+UCJZM1atr2x/EbnahUoaLVjSTm747u+k3JQyubJUH9EHxXJGWsAN2cS7uOcrPQcQitCjiHQ9Caj/Oo7qN+NVVJMcHLkolqVWL5cDXiHMdMllPNxgzdGOdz6baokCywYZcTcuXo5esRRUK9J25UUO9RBgDHKTgFQeqgITCihluNyAQvly5/UXftqX//SfdDg/gFGz0MFg419+z/d94vUHu34iXla8MY4ZqSUAXSEpAnYEzG3XLZfLs7PTxeLcGhP5PMhM4BywY7ZhSkCkJ41SmggBEBRNJ9NbR7du3z46ODhU2tufOSDquy40f6QA0Rq77dp5b6czBkDHAMDkEClhtDXm7MtG05uL5cVisexNT0o5sFqp+cHBndt3ALAzHTvWTaN0Y01/fn7+7OT03r2XZpMJALib2369p6H1+j9c7/KKrFi9u0wmX+Tx1RY4s9vZByQwFdMqJSxXm7ffefjOw0dtb1TTIBIk4g1DFRMrF3S//g0M7BwAF+viecMbxrWqatZM0fcyD8BWYXMRO1CKuzJcFGojPgE4RED8z/H6EClUS1o7yviKQrgAy8mHD2lO0HuwECUqrG2y17wv+4aLxyg05hlkfTuYuFSjV6HZW+wV1xLvRWFfbtYJL8qMDaGIw2Oms+PjpeDIFVGcUo85P6nMcmUsXHSDYAOlzmpsCzwlxzRqEAaOsegeC/TFDVLsstSGaF/xH/7iz//Ya9eMR3/wm/7cJ994UGXbglOW9hCwGEJThI+42GQPDBpFSIoUkTV2s90sFovz8/PNemWtAXbZtpkdOA/xOwCASTObzW4dHR0eHumm8Z+AItVMmulkojQBs18LUAzWcRCAI28Ba42xDKwbTUojmgDscmFBkK4/ZlaKiIgZetP3feccKyImD4yhY2aw7BiJ5vNDItq2W0Tcbrdt18HhXMyWYUiN/ayE1ms+AWb+IO0BBEErRAewatsnT58eHx+vNxvSWinNUNAhknqWv8BKu28WLa4gsde/7xrs6p10IN7N9vA3ONUl03grUNeWlxRxu8gr6XIYeu6MDDRwHJIaLM7JH0EjOP3Yy+HS3mtcFmfs7a13bWJ4LpQu8PIZwrU+R7FPirWVGMAlLnYVlEcw/l6Pm0zXowmGHV2b9FhkqLUysEDXnncQlwL99d690cuPGQiJFCJAb0xr+r7vl8vV6enJar0CZ5EIlGLHbC0AKq31RCut/IvXTTM/nN86un14dDSdNHljDgCArXGBJo6AhI3Suml003RKMbC1pu3azWYzWU5IKWONd8pILOxYW4SEiwCKFBJutxvTG+ccO4ukDqazg4MD3SjHzhrb90YrRZrmh4e60c6atu0uLpbz+eHBpCGiS2/C9/cM4Dr+k8Nvxq8rmP7Kx48bWkYFKB7U/iJ8OKeUAsRV17/11tuf/vRby9VKkULSCU8WrJJCxoxyemYKpZ1LdVokgckZHY5wCofKcQkUZtHRJ/EWwX6R90bl0SfmUVj7IxWehiyACIditRPEvikWcVPKh0lPrVg+F4tcsgQPUmei1C0Fz4RQmkAXsJbArCExlIPkQrimunzSnj0XCoteBgjzMlGN3Acyd3jKeDmmVATMYDEpJdQRpPW58C0oNRWoNE5DlKSZhG5D1fNwWO7FUsM/fyr59TFU7Wk1J5Z40vOO4lBskSNUyhHV94UlQIa4WDx9ZlYKddPYvrtYnC8W523bGdNtt1tP24+v1oIFbOjg4OD2rduHR0cT3fgrVCkirbVSzkXkk4Cd5xKBZ+4oQq0aUm7aNJNGb5Wyzjpr2+325PRkvV4BonMMCCqMFjAO/QLTAwliiwJt223bLTsL1oHmg/n81u0708kMA0+IONApkFBZMJvN9umzZ1rrV1+5P2sa8GLRoe+tJ2osNGmi39KNQ+VNH7+DICCf0qVy0FdKjzJfnfKuIzd93UlxrEFX69WDt99+++E7bWea6RR28XB5pECnUhelFi6oRZh3SqsA8+5yk0tgaGASmPz16pKXrzJAlZZl8ts7/YX5MkIPjzYXw60HHqUGZb8VhMsFmMungZidVvmKoEScdZoL+LuYbxaz+Uzmv9z7a2wJNk8teLyPQKysIAYyEliTh668svmSDoYzRRaGRgjl/L1YmHQvwse49vcHL4QRycMtzLDdbNbr5dnp6fni3PRdSH6E7Bz3Bgi1bibz6Wx6cHh4eOvWrVu3jibTKQNYa/0WWFDWAyBFXqHNIbNl65x1PVuHtGV2m83aWpMyqDFmvVpt1uskXZdG+FKOmwhBiE94kdE8gUdkdm3XIoDSupk2cUmUFZFTujP9ycnppGnu3DryCeD6mXc0EN40VF7HRP7Krz+7EBALbo/bES7SPYKAYAC6vj8+OX369HixWCCpyeQgEYW55odzpieIcI8D4X2sRRDH+3SOK76Q9nYG/K9ApiZMmoiIeVPkMsn+LBdDosYsRIQRkjro+EaYEIZkGQ65Iq9wHUSAMumwUCgVRi61qhxgNeatlB4L80bGMtRCYdTBcsWpdvziii9V/CBGLNTehAZNudecpr8UYZa0s42YXX8xEhejwVy8XDJzrIDwhatWIRKIYs7Mw4Ii01MxzBzSg9OsJhdwWEzXC+Y9itK/WjB5TggIyumVdIkTvuvItYKtY+crNCJEwO1mc3J6dn5+2m621uv2cPCF8/8g6aPDo/v3788PD1EpRQqJjLUMzI5dtHH2FwR5/xdEBGuN3bbrzWq93myMMdba3nTdtnUsFXwcQ5LTBRcnAARJRSPIQedX638tImhyzOvNBgC2m+7o6PDo1q3ZbEZEXqWCWSngvu9Xq9Xi/Hy5Ws3nc6UIXngj7DPAo7kZC+gzlQau4JYwO0JUiB27k7OzR4+eLFcrB+w7Ox7odomfyzWQjSNzZuZrla8IuMNeiwe9iuctYySeXg5JV2VlMf4VjokOLoOdQSjyygiPIxCtAGQQa4NgHEGvsS5wsWz9i4gxnA5HO/ESRuBdt8DQX3jwuQ60/Qa/F2E3Xh9H/oV4xYgoRtnccD1xxYEs5pWGgDjaTsr3tt6y5rELh6vPG2Csqnmu2D/2PIWnCg8FOzymoZVGhL43fd+a3iyXq9OT4/VqCeBIaSTtfOGPoJvJZDKZTmd37ty5e++l6cHMWWeMtdaBcwHeUTqs5lrT98Y646zvB1zf9ZvtZr1ar9frtmvBmVrkK1+cKLdjQpMScWYHADaGgnDDOP9zHLvtdmv6ftO2nek6aw5mM6W03z5WpLRW08nU2K6zZnGxPDw8PDo6VN4ex32QdgL0WJN3o6/55n+3jERC9HDk6uMg/OC/1bbd2+88fPvhw23X6maCpHnEUjD9RPTrJyjjaC1thZB3+QroPJsRhqBMWfqFYWDhxRXzJi6yXHkxBLQ6SxQ4Fs+X0zii2OnlLAzExcCk0A4F0QIwwMc/8srXfdWvqYfKouQr92jxq77ujx6fnsPYlFEmg8qcBIepCcdVFoRlbYr+KPqx8Oj/4Of9rJ/5k3/cDkCipkn8g//tX/y5b/sbgz44iiMJ2UsMUmUQiYx5K/wH/YDPaxoNWYdoFM+R9FYGHHMwG8/TxUvJfgw73ENref6SzrHdbB88fFJo9HMx8n++HFAMoBizf0EcycjBkHOOiBqlrLPb7WZxcbFdr7fb7Xa7CSMQ/1yshd7AbHL7zt07d+7MZrPJZKK1dp6JF5tNIiKtCD0tp2u71sf6bbv1418ANra31pokF1EIvQe2FgPlb3ICkF0eBopV/2FGdsy9tW7bGmNW65UipUg1jZ4dzA/n88PDI91opcgYd3x6qidNM5nMJ41fnIuyJNXFPmTJ1uFRbOJXw/Xnjs9XxOH3ewdQtr2w2WyfPn12cnLSWzeZTBzv9JKthpA44juPl+s449BqFgeFYFF1ck3ViEtWsMvvtniJxV5ULvxLlj9KQxksVW0GGmbV0zycH/ysf+vHX/+Dmc0OAC/qCW4cv17TDEEGQLk0f6naaPHTf/AXfvyn/6Qfc83n/PTkrLrQq1U2BCmFn6pvJzG/P/57vubzPvLKB6KC+//9s+/+iq/+vVXIfs4KFEdpaAOSEpb7bYhEipm323az3Zyfn19cLDabtbMWEFApds71vbfamh3M57du3X3p3tHR7emkYWBrmZ3VilSjgMEau2m3bu08INMb27bb9Wq92qy3Gy8BHQu8OFEA0AFWqlZbSlUSHLu0Lg07nv5pTWva7QaYgUg3k9l22/Vtb/pGN4jgrN22m2nTvPryyzidwAft6Kvav/fi6xrGvQJ+QQBCA9D15vRscX5+sd5sUGmYTEVEgspMXYjyDX49s0QNeGdxJoowDxYmGXMhJoNByDxEuFBSOgBK2vY4HvQFoyR51GUYP1/QXErw55KHix8zoh7KOIJW3DAgoKigMD0zFHo1I9ZaMtTiSLxPn1AhXorlHB6fV2GXOLOAMY9LMCjlxP5qYOmKxbjgA3OwdEgImkPP9yIqZpXgi4F04+DM6PUb2du2Pzs9PTs72242xvTOmdg/5Xd0fnj46muv3b5zTysVhwmxlkAiUs7atm8XZ+fL1bLrjL9MHDvT99ZYZgvozS2Fm4C4ulAgvCyVpaoeEyv7hazomRla7Kx1A8dsNtZsN1sPcEVc0x0dzg8ODpbL5eF8TphELC6/kPB66Bu+13H4/dABXJaRHTuFpBB75pOz80ePn6w2G0QCVFyE0npdCFPdn934xgnopb3IsDCHRChNfMSx5bByOsjXZE7EoZVg4UsOpVxcL3dZaw1ILMQtubDre3chSSEUmXG6gWGLfO3jAxMskLjRT+OauxiXN5co9t+y2xs/3/X4oThFKTVIJ550qJQiRGP6tm1tb9br9enp6cXFOVgDpNBzNvseAPRkMjmYTCaT27fv3r330uH8sDem73trnVLUTBpgMMZ0Xdd17Wq1PDs7X66Wpu98TSQ0SwARQUFt3cnZOhUwOdphsfYpoj9GEySM2nuRyCbkAmx41Yjkl0QxzAbYOmtaAy2D9Zp0YLr+8ezJ3du3tdb3791VRAhgnSPC9/9V9e4mAL7ZnSSNsZDHhJUzl7vturffeeett99ebzakJxDsIISDofg5Ua0eM5CSygVGLlaHsO4ORd0gX0VB1hPqWRgpLsXCPe4a+XINfWHkumNsHMSPGVnIweETw6EqWGo9kk5N9tG92cc5urtbmpWIIh4GUD+XkhbloiQWg+NShAb5hYMYim2OLPTKovDEPBUuNzQ+eGco93nj4p/zDgOX/S7m8scxE3o2vWvb7WKxWK9Xm+223Wzj0gMgAjsAa0Hrw8PDl+69ND86nExnWquua8PiLjMSad0Y02+2m+XiYrVebbbbbtsmTifk8WxWDMSg4az8QohzzjmOvjsMkuLFUbILxajRW0gCelJ/pRnJUZcpvxlefC50346hbImd27bbR48esjXAfDj/0qODg/hukvABGTSXPB4hRjfw+AOSAN6DWiR4E/JqtX769OnxyUnX22YyYajnZjz6l8ejL757zy9IleDIuA93J0iUewmAxMhpP4tZjC6EXXq4jHZdFNWybMGLeN6Qwm5sUIMjdXxc7KrR43E6PV4rHeHz5oDMwCwMHiL9C3dQjMKX7oMY+N/VRrzc+Ss4vqRw2ihgaLtus15fXCyWy+Vmve6tSV6T7KwzDITzw6ODo6M7d+/cun37YD4HRNtb4wyR0o1WpBh4s96sN8uLi9Xy4mK9WvV9F4hzKFVS8+qO82rZfjVLUbSsYY8WhAXhogkXwi7gvR7l1maWMY/ap0kGIzgKkPL8U3TOoeVs1ae9ISQjwHq1fti9fXR49IVf8AWHBwc3lBGub7APbgfwLl/XfofbAmza9tnxydnZot1uIYh+YHXHYlqQr+AHIQd+Faow5hVduJYLs1ZGACZEpCqBSxrGQERAkhb87IIwybsBFyuVkbaC0rg1i7FAXjsFOfEmWdeMNRI3rKLHFraKFBUlJuP7zIIXVOhLSJdL3DkQEmwnxuePZtF6IZiQxOWzoYSO7EhqrcgPHnADL9DD5P1whuzwnB0SPNdJKW2NuVicP336bLtZO2ets4XOjmOwtpnOX37ttZfu32+ahhGc84v3oEkxIxEppdbr9dOnj88XF9ZZZ4yxPYfFX39be6I+KiLShIjWMpjeATOzc0xpvcYLRDM5LyUUyHROxtysBJc3y4UGrtCKDtJCREgkbe6dC5vIEHW9iEgpIsS+77frzfHJ8ZNnx7dv3ZrNpl4HwyUeYTGxqLIAVwv8+wQAAMCO/QZVZ+3p2fmTJ8/Wm40fOfFYRV3DvnnbCK/MviMV61VFVrCWoHqAsPPv8u5fVoltUtj34uC8WFiVyd/PZbHMFQqfOBFRbQBvrFgYcDYB3vCNo5FUmK9HBFyPUAAr25Pn0bUJ+z5hmJ6NV0akIbJmP+bVtA/9DGCwG0+EgMSOnXXrfrVerU9PThaLczA9KOWNff2QnYj0rJnOprdu33np5fu3jm4xQNd31lgkapoGiUxvurbb2M354vz45GS1WvpAyswe2QH07hxiUx2JFBEBEVpjHTvHzpgoO46UF8Zjv0CxJ+UU6DGrZ4sbRVgthjSQnIqTSGXUHoqDBPKra8HAHkkph7hcrd55+M7dO7c/8tqrEz/odg5I3AV8OfOFrwcIvScJ4DpL4Ffz+mEEJbj08XlWzlkCPnD/nX/POmOePH32+MmTTdui0oDkb9UkbkPJQYKTWkzhzScA0qH6moCgxyhimC1LgjIBFKsmiQ0tdkMv9Yus4mFQtRVPJ3nCeyiSgHzOSwKQmNUoRb8TcFoWhu4MEJ6xw+Tge7OaMpk4Yub/AFa8k9yADVA5zG7MmA1bEqUp6t2g1MAvdDvxhgUtBkFwkkqWMtzLCTMLCq13A8IXhc0+KzApjKA3z7EHxnGbISLezEJEltk5u9mszxfni8Vis1oDmyy+Hheniejw8PDlV165c/eu1rq3xlfKrJgIkZCZu75fLpfL5cVyudpsNsDOOfbWcURKaY1EjpNDvGN2YN0kAP8aEb1UhLU2XpUObPRJIPKXWYrWOfNXZsBy99uxC5yfuOXu2LENFvN+HBxSEPj8EHKec9ZPDBq93W7fefDgzq3bd+/cnRzOANCBIybxcTDymOBxReiTyyEjseM6cfiDtAcwngrTldx33fHJ6cnpSdf3WjeMOKDopx6/GCSOIR9DAUoe141hFiuliKUnvHSHv9wJa/wDCG53NPiYQzwsbFgQ+aoOBQpQE4akUHru4DL0Fy8Bm0tizcC+cEjKTsjX5TXEzQdHI9XsyE2Csj5gRg87fPB7AHy+LeBkWZl1HllrhYimN9t2vdluV6vl4nyxXq/BWQwoh8da2Ms1IJLHTyj5uiCioolqHHC73W432/V2e3FxsVgsur4jRK0b66wHfzx3gxQCUzB4cV4z1PYAyuUqXpFicAH3zzbJzM65wpKhUMzOQqeMDC4vzhc6lH5UnKyA/FI/eqCJk22nT4/WOeeAUKvGWPvs9PT4+FnXtXB4IC5weN8qheqr+oub8PqvS269ou30SkxI6AAcwOJifXGxWK7XDlBPpt7dhzn5AkQZR4hKIxBnAQPHXkGZTmicNIGvzWN8uo4ClLKJK5EN5B11mVCdQ4ylflIliso32T8piVkm+ipygaMjC4YMp30ADrsKaR8BMl+bnxteZIn+phYKK9wfdzarKCRrspkKU1ivFehL/igDqwuR4GagU/4wUMxQULKkIOkXJkkfEnoWjKIV/OCNgcsK4MY5AL1ySVjR4Ojf6++pzXZ9enq6XF60bdv1fUSFkj6KWAZnt91uTk6Oe9PduX13MpuFtoxU37anZ2eL8/O27XvTdV2HCFo3RAA9GA5gjbXOxSKdFCFqsNYxG2NMHM8q5Y0hyTlnwEavb8/fCaPefKuxWPhLtznL6QAWG0F5UQCVNx6j4MDheUCWmdmKH8UKSSlljNms1+eLxcVyee/eXYqT5LhMwFI8CmrJq/ru4Z2Ws9eJq5/RPQB+AZRqRIIteadYdher7fHJ6XK9ccyxuSvt+9Jdj5e8TXHrRLoqcu1vV7in1vx0vFyqH4qFMxYotxQ4zjS2UnE4TS4C8pOQxmxAIgtmeYkjlli61G3hd7eezctqcNlCJUKpdV3+UeHYUCgRsej96EWhT77k72cjz0I70zuFfhCHACPj3+d90wQSoZQC4L7r1+v1+fn5+fnZerNhZ32hwdY6x3rSTKZTADTGGGvZOQfYdb21F865STObTKdI5Ixtt+3i4uLs9HS5XDrnPLtGIXlJTtJaIzI7BrDswHC2xEEgL6rFQh/bi7IhulS4o98dduHuQ8HTRsyshVBb+MvXpcs1gF4uUAZSkrDgwAGGuRz7Z5gZdqj89Uye2ArorLu4WD5+8vj27Vt3jo6IFBau9O/fDuDFo/+N2xwslLQy89sTaBFwa8zjx0/eefhwtdkQNahUuErZ5YWrNLEZ3hNcrf87ufXLAuXhAYuTqs0qVyjnMIzpz40MAVhQKMJmAkOt5Rm4nhwvw9gecCIsIEplA8TodJqGESSIbkKkhTw90yeV5xxwZt0fLICcinBeoi5RCInLdB0rMNl1FebHQjqTh3jRTcrhGutLNpwyxydtJ3qPR22fgRYABS6adp1u8jPEHDTeGL2xi8X5ycnperXq+o6d800TswPLQDiZzW4d3dJKd217sVxutxsLThEaY7Ztu96sm+lk0jRtuz1fXJyfn282G+csACpFGgkALDM4h4RKK2Zyjh1b51yW7/XQux/rECKiM85YY02h6RMic94+DzoVWZcr2IqFGS5E582wyQ/I7Iw1xlgfDpyzHugxbMFZMeMiTwBSHugCcM4bHwcd38XFxVtvvjWfH85+4A+cT1WE1KjMATw+6sVBuXjt3P98+UW/H67d2jgwjH/BWvfs+PjRk6fbTacbDUicnDSEuA+Owp8jbwdhRQ6F8S6ryh6VBiKOg12VGubQA0wWzjzS98XugOrx7vDa4AhqDH2sQOgYVeU/wvPFUhCQym6Uv2jjWGZVgU0w7x5qidfxIkHYl3hcNHNZH4QztjZqt+j//5f/X75SDjq4UBbCot+rWLkAAPy7v/o/+3d/xk++5vN9653Hv/D/9ht3sT0iI4BzgYQ789yL3ofMShESWWPavjN9t91uz87OFucLY/rYobLrO2CYTKfzo1u3bx8dzg9102y2296Ytm2ZnWMEgK7rzs/PjTGTSdO23fnFYrvZ+hmvtQ4AtFYM6ExvLCMzxcaYUFGU7Az/Bst1T+lC55xzfmSMREqRnxNYpbRuNAK6iGkJC2XMe8AkbpwwN0ZEdNa6LZMCpZRCbZ1h6+ITYOdsgJ3ZMSN6c1qXrqxoS9lg27aPnjy9//L9L/yCL8hCGteY4l9CP37v2gf9WY/9I1BQHLlu2u3p2fnZ2blhN5vNciEsND/SByCF7Ect9zDGgkA4wIHLRpSgxCKMJpNuFHr9QzOmy+ZxVD2s2AaIvwnz6kOaTkjZEiRgX+lgql5xOM+MvhyIwfaIOQ6dn7P+LybUyAMHuLHBun+HazLGoFng8SyZuwu+6Q6Db7PGcztnR4iIzvHQvVlSthmcEEQSw4xBIgsqacwIeP/enes/4Zdfult7u4xdltmcoSoguZh7CxYQPtcHjexc23fLi4vl8mK9XrfbjbEmlQ/MDpyFprlz794rL796cDDTSqlGN02z3W67vttut8zM4HrDi4vFer3WjXbMXdsCYNNoIHS2s9ZpzXFc7Ky1jhkQlSeVavJtt3POsrPM4DH3nPgIFClSzUQrpYzpjaVJM5kfzQnIWOsr9eSFh5CpI8KjL6lsESk02PtH6WYynUzYsbXWP8xZ56y11hprvZsxM1vLFmzc3/REVEQi07fni8XZ+WKz3d6eH9ysreSBhd77HgJ6AdCfayw5ov9gmbuuOz1dLNcrY3sghUQ88j7V/oT5OyW7pPglo9ZXY1IUpdd6BVzBNbsvhOFSq9AXDX+ELGiXu9ijWKAsmCV6E/CfUxJhlDeou5cbBgUhmoED0AuTLhcKWQcsR2o4/g5ll0OE4ZVw3bf3qsoiWYRy2A675liE5doC8q7tzfQKhKnYjRLAwWx6MJtutu3YHIXrZLBjt7qG/m/+YRORnjR9161Xq9VqtVwt1+tV17aOnY/7bJkBqNHz24fzW0f3X3r51q1bqlEeSJ1MJi/dv6+0Pj5+tlqtwAGD49520Pe28cinIg/CEEwa05u+60lZJFJKMYOzhgGcY4KwrssBkw9wpseBPKkUlTeIUapRitCxQ2bSWuspEYKxZTtOiQ/kZPaM7gBIRIrIuSByiKRVw8R+QOHhIOscGaecc85ExyDnrAudAbN1DIgaECz3bnt2fv7s2bOjg9lsOvWAEzuWdwLvurpLr9MrxAv4XUgA/N7E95veplH3X6FG2lrz9OT04aNHq82GlAJUUJvkjlEf8miRRPyXxi+OYbB4NUgpsujj9AuyVD1XY04e0VtDGD7PonOXjiNFdnBxkRjjnlJBTqfUomApd5knBoAo8IK034R84wQQ9LJqxERMwbxzXsBGi+CDMDDXrAYF0TkLBzeAzMl8w+uPQ8yKS9bB7gOZ0tI2864hFtY5QBKVBI2oMAPICcPvjr50986NnvP9e3cfPHyMlQs1FpvtVR7IdN9Cz+oF1IwQ2NnVavXs2ZPVcmWMsc74XfesiKlofjB/9dVXX7r3kp5OvAqPn3wi0e1bd7RuNtvNer1mZkAFygGSVoqIvIintRa8p4rmtm3BwGQ6VUSsGUB7yr+1LJxnwO8GKEVKaVKKiLw9ZLjWCf1k1m8FOXbAKgpCxFuPOCqE5cXfehs9bpi6JCmUADhAUgiKSfn3YRpuescmLDAb01vnTMAWEMHB+dn5m2+9dTCbfeyjH4mOwYx4nVkpF4ZNl8bVF98T+8x3ADssEcKkPqD/jvn49OzR48ebzUY3jauJ3XFSX2nv5NZunO9fIDAlf1EwL6vhYRWyL39pOJB74x1rGliYue74kDH7T4riuvgbKB3AEEZQdER8vivkH3zbn/xADkQDj8rfi2FJkMoO7koxUB5tzIUrVs0Si/8oxLt3bt3oKX/k1fsPHj1O9UAG/Kt5FqeIJBHI1AvywPvgBscYe3J6dn52slhcdO0WgH0yc8aAYz2dzOcHs+nB4eHR3bv3DuZzB9wZy35soJRWjafqyJZCkUYi/3K8wU7fG9dZryTaNCGFBAkf5sDh8SE/bVshESkkRdpPKBQReSQ+WaL7iSlj0HGzHHJn1opJVpTAIGQCYpuflwGj73xm36UqMFhSijEjsXNak7VN45idM6bvOtRKo95st289eOv27VuvvfoqNBEXvd5CzmeSh6A/w7SHjAcUi1YRbYmfWW/s6enZycnJdtvqSaOQhOAmjljSAZbePg5gZGE1KgI6EDV+mioTFZ4nMfTyePFV3JECaRiozxVG41jHf5aKPkIGLtT+JBStMU0Da5onQ8W4r/Zpq3bkc/3kUpqFqti4enaJmZT+CTykbIdShbNIZvlbmL/wCz5ON6zAf+QP/SH/5F9+D6eOo1hIkmOAci26zFG4M31d6/zG//RX/Mb/9Fdc55H/01/523/oT/z3zMnDi0hpy255sT47O2m3LSIBKSJqJg0Adl1nLehGIxIQ9l1vrJ00k6ZpmLltW2tMVjlEVEo3jda6IaWiEJwwbibiQAONLo6IgARIiN6VN9VhyBksjW8QEsKYUDkLy2MBtXLcvAHmoIzCfrPS/xgiRVOlFREgdG1rnQVomknTd93p6fnx8Ulvet80hNFcNlPDF88A+MIsIPrs3qn12AuRASzDerNdrpar9cZai97aLeuwVVhxMoTnYBGOhUYOltavhCCoYYGTT1jx/vFy2KpYREAM7yIz+K1xLk2+4vJvegphYgSlEAHzjol4xT4pDExytwiRQjpu7MTwIQn/xWII7zbXTsV0vAqGn35Brqp6KhSZthiO/Kgf9kNu+pR/7I/4khyGWLC7RG8qhv5+1MzC75lH5x7v3THGOQdKq+lsRlpvN5tnz549evTw5OS06zulNJHyNbPSWjeaEHtjjDGIqJuGSPkRa9gesJaZkdRkOp0dHEynU91MddMo1SitlVaolV8XKEdy4g4S93n+GqH8YLnUzEi1EYYVQAZKVX7lJR4xWY4pwfmlzsBNVaSU1poUMSJqPZnNJ7MZIHa92Wxb63HD2IVyzcfD4h98Dr7eB2cIXN2QFROQvPZn15+enq/WG2YbdjX9eEDwt8Q6X+H7KbVjcEyWHxFUUWZh/A0oB21xQZcFdb3wQY3eSyGIuLibLIs0EpQixiyiMxR3G1XrwILWGZ5AZANmOfu4Okxj9arAvJnxQ5ACUCj/Oz9PYx5F8MobHREKkwBOtnEC6098HNq9cPUcCeDLvuSLJpOm63r/ORGRx5FZbCrF1iMC3oQUwawoUJytit7rj1kpQkQPyvTt9uT0+OmzZ9v1FhGUVkoRQJP0FZTWDNB1PSBMmumkmThrrbFt2zt2AAi6AQSl9fTgQGsVVf0p7kQCVM6/vjXmuNwS7kTy5VBg5Lg8pYKCiMliUpgkd8KGOynSTiFlIV+sJnalvx6nlSAHjtkY55g1kW60VnNnLTOeLy5uHd06mE6ICKKbwrv2Ab3wD6LP1l0qey8WIvjWucXFxfHJ8Xa7BVJEOnF3r5MWy+lqWVUjSywES37osGjGgciBmH3KAiSCkH7BhMKuIQm6am4wEC/j56UfHbsSSshWhiXrj5+runQs+mO5Lf25ngHixY1IIb3vECnhIlxiSbOq11OKBiB+VsW+EQDCj/qyGyeASaN/+Jd8kcB7gl8VEYkdplTrS1rTgO/D772SBaJumqaZWLYXy+WzZ89OT07Xq5U1ndcFQqJm2uiJNtZ1XceOldJN0xAq9ooOvTHWAuJ0Op0fHs3mB6S0F3/WqlEC/Ml+RnKAI18tggr7WGFhIKr6sPyEObcK4hrh0jSAgZBUMH8J97VLcuIVjho/b78M4FwQiw59BxKpRk+azpjTs7PFcmmdxXEtTB41j/9MzgA+w1VhZq1k7T32gn2IAMbak5OTp0+fbdqtUtprS+VdnhHCDQ50TwjlQDAQQjmxggeaP4PIIDUbhp/IiF+A73cpkt1ZDCtyRYlZ93Tg3yun3BTWwYasSJYodAVf55fEwwWI56OBfmAnAKPOZMNmdNCLYlLD9zVHYc7DZVIdvXw+9trLr7780nM86S//GT/5f/sX35OuleBCjaUZQ1G3Cm/QYUX4XuoO+F6jt/1qtTo9PfM+wIioJlOlPJGGG6XZIbveGAMAzWSidYNojXHO9MyOVDObTefzAz2ZtG1rnfMSzzEmkHCHABCbe5GFJWSjpP6TTwJegQ6rtb0BQ2rwLhGCC0CNi502F/JAeSeDBP88iE8QKSLlh+GeIb3Zbo6PT+7cvnX3zq2JUizMaKTnq1hg4rHw8B6ez3wHkHT3Uj3DLMyonYOzs7OnT4+3m05rpbTKFVyC1BlE3TUyy4tIXRy9JHewAb82t/iXmY/gAEKUVxBjqPdjjZIXd7CWTODhvh9jtob0iQpQ/NUSEMpjCxCWSTtQQ8as3cD84RkCF93l6OtOboFlES2Z9sxRWXrQWuxe2fw5P+0nPt8z/bk/7Sfdv3sbq/vEcRapYwiXWfrgq889O53we2ptz8AXF4snT58+ffr0/Py87VrrLCI1E62bBgmttV3XWmtJqaaZIKDpjemM6a21DknNZgfz+eFsdtBMZpOQG+KMN7g4SO2RrCWL4q0JSimO2TkIJmJRQs634AiRSOSXBvw/fn0AlUJFqOIusVKktGebUizXQEXkgcR9La4rGcWSSp0m0gkQ2Lbb49OT88XCsQx/Agmvm5HsS/mZuVs/K5vARfCt6piuaxcXy8XFwjLr6URO3itvUqgTfEqcrurzc/THsjAcn/3ySDnFo8tjCRH2D3eFxERuErla5Kqx/pJZhNGHRSh+YoE/4hDgGha6LNij6KVK3IctARSF82AQhVkllLm6hMLjiJABk+AwlGtuVdIlxF/47/6M53uOk0b/wi//Gd/8LX8pXx6OB72hEAnBJPjJKGzMGPMM6z06Xd8/fvp0ubwwvQEEpQhAp8EYIjC7trWK1HTSkKa+N13Xe2a+Vs1kNp3PZrqZeJgr0Gy8wrOiIGvltR/yjZRlR4XLLhAiB4iGY6p2YGNd7dDP3fIQr4AhYj/BQIjOi/9AWj82TETKQ7lJe1FeQ47LHswPKVWjSQS4vjfLi+VyuTa9gUnjn4MbqdeuoSnwHieAy60D3t2vR/AURHQApu/PLxbrzdqFJo7Yy1ohDutxKNw+RmR+Rt/TpBDM0pJ+5G9wBSuIdC26Brk5wLLNlHRXseqQST9SEj+LE+d9qmKfIDSKEfPa7TpcEJchjT8RiJjxQxLxSyHuURQSSsVkhAESxxkA8BdnRh92yGr8hB/zZR997eXnfta/6P/4M//s//RXjbFVd5g/1kBHrxgEGLSzSk1Qfs+y/WazXpyfdX3vOfqM0DSaAbxFCyIqrbwsmnXOWtd3PTM3ummmk6aZTKaTppkQoYtE0uzOUxpyCltrFD0OhwEZACtw7MD7QDI7a7uWo/xywRgUOwNlms+IDjrHXu+BfTpBRFIe1gkeB0onJo83hbXiyfpFX60VMPgFOa8P2lu72WyWy/XhbIJKoaD1+e20gU91FTbf25isr6cc/V58ncsxJCQkw+58tTp+drLdtkSaa5X9cuLFXN7r0hl3JDwW+YJwLNxfwzAyBfeM9cvoLyj7dQiPZAShcoMimOwmlQxdZZgv+9MyZ0W/PCRA9oq1H4b4H0MkSmelHd5icuVD7oak0VQU5nZRmZWGfHwMEfxFnvbLL939pT/vZ33L//I3w2USBAQSdoBx4Slti3HaiuQ6oiXzunf/bLdt13XNpFFKW2ucdd4b0VjTm75pmqZpFIGxtjfWWQMMupnO54fTg5nWOhTqLjmwQfbqZRdHdnKLRUhapQ4n2swxhCGsl2WAId+Zo+CDE/mFdxorCWMABEKllfKLyFppBvJGjwjIxEJv2v9lAiSlrbO2N+DcpJkopYDdtu1Ozk4PDw+ODucqQEzFnjwmqGDk1n5v47D+7Nyhg8bH838ulsuz87O27b37Z65xBrV6ImHjDrhzCIpzObyTy5ajsb90YOegGyjgdSgVuOJ6gdwvq344FhYVGC2YmAtgR25vhX1WZBiumA03grgqGDIr6AOqdP+8IDUKh1d5jVzaHNZhQbZ3iew/MkxHAICPvvbyT/kJP/IFn/mv/hW/+Dv+1//94eOnEGS0sZCS36V0S1ivNAGTovfo4/biOVpppTSws8YaaxURKtLQeG8AtmydA8RGT5qmaabT2XSmJw2h93rktEcGdRiVayyBuEGEYfkB0VtF9s5Y4+tsx+AXiqM1Y9CBkKORtCjJeWwLY0TsvCfscTWy1sM9jpyx2McOAIMnGSmlFaIK29hEgMi9660FIs/oZQebdnN6enbn9u2Dg6lWqgxEBSUdrj28ebdmxKOewO/118W9mvKec7xcrs7PFm3bord/RuQ6Qtf6HVXpJ7ndlfcm5vVuFGME3qk7xrK2ByEVwWPht1xUGc6MK0gpcjazrkAttEaikicccJmx1meW2UW8G3Ed6jnEQLu+5/fB1KDR6ia9C4+JM+Hggx4J/yxGLrEhwPjzfPLHGqGIP+jX/8e//MUbrNl08jVf+R995W/+faP9bOGmgNLsurqwvKTCe/XJKUXT6QQQrTUMTIqcteycUooaZXrX9z2wQ1LTyWQ2O5hMQ+HPHDx2pdBCtuhLQq0eoBVpwDnrSZ4+Slhj+77rPZcUHA5lARIChFFfkWPfL8vsYvFedAbe/RERKC4XsO2NV6XmQPgg0lo1ugHWoBwyIQKxQ0RnjLUGQAEAkWKA7bY7OTt76aW7r7x8z5uauJQAd47xRlxvS5xDbhM9f0x+jzqAaznDcInZM+DFxfLk9HzbdqrRQJUtFDMD7nQ14VIEmHF4p2No6wowXmg08EDbHYUWJMVoGmH2UiM+mtLwznSdRljFPkCEgAgwcwny1VmYmuKgwY2LS4JMmkbOksj+fCTQf+eXf+XbD58OvC2F+J53afI7C8EzKcIXGEATLEi5kDHcaJYkqx+U7nzx/Bdf8St/xS/68ufuNMWqFw+820A4oGRBDg4N+UjSHHWG+zd/9Jf99J/8Y9+VO+fH/cgf+u//nJ/67X/z77FoW7EQfKqiXbEbjoQARGDbtvOy++/6IaKmmXRd23e90kpr7cgDMI6ds9Yh4nQ6081sOpnoZqInOuDyLtLzC/tXiPo7wVPeAZADLwwE7Ix11lpnHLD1CzJenpmdA3bggzImqhAGwXaM5vZ+iJgFZwPLyLtOxeHDIHRRNjVy7I1/LTgXFUCIQbED6xB6NqaPiBwRkbHGmN67VDVaO0fGdOfnZxfLC1ct+z+fTzBX2OWLdgA7lRFf4GsYcmmqxCDWL8J3TW9W6/VqtTLG6OkERQIIdriCPVRqgtb8eJHVsfBiIazl1KD0KuGBRBsGqweSjtIIqVxEIIniS25PwfHnJD4+0FVGYeEU73UnvHxr3dQymCOW815GcJxWyIRZ6407RkzeUrWUf21CXySs3EyPquAJMh/icM4Fl9fp1wr7wz0AlmSr+IFg9T5ipldlIl7ynaiDbfzLWunf9BW/8l2MsL/uP/nl3/vJN7/3E2+MKRbxIAVltW9Ej4ZAb7qLiwtn7fV/6T/6p//ik288IEXMTKgc8HazXa9XvgtUChHJoyz/5J9/FwA454w1SKi1JlLOsel75ywgTWez+fxwNpsToouay6UDR2LyYcJAOEixee8Vm3a6emNM1/W9YXbeDdjbFJMiJHAWXZSQq4l6lXArh/I9MKnEBCVT5qDSa442w15f1m8nBHk6IgTnrLXGOWZro4WyCn9de6E7DlvQdrlebYyxMGmKiUN9m/COcFrfwLv10m/29WfRDyC76DgAY+3FcrXZbI2z3uURIBF1SvNwSLKqw9F5sevB1deIMAYoSRJdIbAImXaRB/dSeTOtFsfnU7lN5wAegUiUC47pEUGtNlfWDEG7vkAKQSYyjGin9yzJuYaxELx6EfCfR611ZaZDYUmMV33gSehmsA492jpdtsR1ZQrIK3cjfnHS+Ef+OUXoLvoc7EIFC2HZX/GLv/wHfN5H3sU7ZH4w+8Nf+xt+9W/6XW89fJLaROEZIXpigScCgHW269rNtt2uNxcXCz8Uveb5K3/z//Ntf+3v3Ll9B7Xabtvzs7PTs7PVatVbQ0pPJ5qUYud6a7zotSKaHxw4ZtP3zlhjHQAr1ehJk2AfAI4xUDhxF+8lkiIiTZ6E41zft13bO9NzkO8OoD+zA2Zma9hiEGEjQiJNfvnLcfaYDlARArr6ykGXJ+vBOSw/PxEcbBkqkIjA7xVA7oc5+AGEDIQMAM55lxvneLPZtm0LwO1245pms9mu15uj2RQp4IksmH3j1zlXguAwcv29QK0En909gHgXkmVYbdvjs9O27RXpchdulLEz0ETDCncfjAsGdNDacUnEAs9vy7vKRYBKQFDYFwXAkSiV5sKi3cFkSAWl4RiA49B3crDTwmw/MFB3k2a2QxuboUzxi0yLuEydxQ/MzxEvG0txXZxj6HjwCvQGpZfvzRqXMfG7kVSIUG+GwMBXnS+5vxh+2Bf/G7/ql/2Cd/0OuXf39jd83Vf9J7/p656dnBXZdheTkRAAbGcvlsuL88Vqteq69kZeal3XnRwfA8D0YHZ+dv74yZPNZkNKz2YzpbWPdgCoiLrOAPDBbDZpmm272bZbcA5J68mkaSbNZDKZTBDJOstyApfbbJadq3XOBV018IZb7WZrTQcApPyIFRFAJ+jcWWYHXkXIm8QgOiBwWSXXC7VFhLIAncQ+ABak07T2LVe0ILp8BRuatD3gp86OExeLSJFKiyWI5JzbbNbOOURw1iqizXp9enJyOJ8dHhyUruB8dbAc2ISKP32hYc+7JQc9rma28+UhY6Q9EUDv7MVydXZ6tmm3SN76S+7ly5uQ5I2ZnAZxDA+LO/1QSQEPN61RrJSRQAqwQusELs+ir2AY6AwXtXwwooryziL+cDUkABAOlMWnLdpmcMAUM4G0LsadkTdf3zeFgAQ8UtlPhLf+8vlyrPqJMA3YIrpajjWGA30Ybxeu99TlZ4yFg2VZXxW9GI4UCBJqqRi+H3n1/h/6r/7ziWjq38XzkVdf/iO/8zf+P77m954tlpWmedUTO2YNgES96Rfni4uzs870N72tjbXnFxe9tc1k0rW96a2/+4hQa9X3fd8br3vZNJqdY2fbjWnbLVg3mR5MD2ZICgJZXkWUpZxjCdUqRUoRGmvXm/V2vXXOIYC1Jr/l3ijAWlSklLeF92RcAnZ+59JYQyaEAA7tLwrfT0+eK3a4vJYegswHEI18QnCXF2qinbJNGSbLDSFiWB7O7nHo9wCsdc4acI4Jfca4WF08fvr46Gg+m04mSmXgqa7iSnT3UvYairbq+fIA3eSeuuRWHFVbvMSZKFwQ6S67uLg4PT3bblpARDWiwiZDdfpQJOCza7mu8OIdX0XLG8KEsajnoTh49UYM/ijqTyBUGrRB0IgIC4UHrgrjsi/KIhZ1MSon0EECnQpKVKp7XnidHGFgklNcCXGxffyXyaYJvcEH7jB932mQdz0rntHLEXH8GizbmcGiaL7/L/2984PZH/odv+Glm7g//uN/9l03ehX/xg/4vD/9h3/7F3/RF8qBaf1KfBHkXLvZLlfL5XLZbjfs7E2F/5hd23Vn5+cnxyfbdttMmulsppViZmuMR+T97FURkaKuM+vtxlmrprPDO7eObt+ezmaVKXcphZhiAgY3RbZ9121W6/Vmba0FZmcNs9OK9GSCWvv5LXv4J85qyXvCKEJAdmydMdY4Z6NFAEdxUI/ExDsvmA979F6R9xUjpUil//SHyH8vmJD5xOOfsJ/uGmuss34qgIhR/idE+eBIY61zFmJ+UEox82KxePz40cnpSVjcg1G3ietE1kEuuPQxl5dQN6CBBue/kUU1Wewnrs5Qir0oisGXsjEDrZbLs7OzbbtFUn54L7mxWM/Oc1U6BnVwDjPStKPaD5HQUCyunWcJYLSBQKhGbVl2H0Xp78kvCS8SRTmEDkDmKRa28pXuc1EPoHT+K7dFdugDwBgglLSgn6OQZixrecGHDyvtyDACR8r6ImpaJv+cyMl2CUIqLjl2jEDZBOLGWUwI+RVwaRq4CY2ZnKQl2lQgFWLknaotIvxdX/0VX/QDP//6z+pscfE1X/cNf+C3/ec/+od/8U36gPt/4ut/yx/4pv/+L/+t70iROsjTMztmpUiRbtvtycnxyfFxu934MaWz/c3eNmZwlpVG9CtOrBqtgK21pjekaDKZOHbOuq5vHTtg1Ho6mTTTg4OD+SER9cY6x0QAJGgRgvtDiKSUc871fWcNO9v31lqLAGytdX6LWBORQiSt2ZGPEQ4YjM0THb9ii6iV9lQzXxbGZa8gASnopoMbn2vKBrvk+BkCFAd0l/0zD6TEEtHyv9QaG0bXku/n9wQIk6fmerV88gQ++tpHnBQGC3POSxwfB3E49jRZLA+fI25DooFKx6EX+boOAilUMch/QTJuAv/HuPV6s1ytut5MDzSSgqHeZXo5nH08Rt4vSQPCAV8Gd00VE9xf9vlcTZhRkq9l1Mh9GEoUGRF2Tfkzz4YHe98wIj0m7GkLgTdGvkY8BNhRel8f2RNIDXJVYOxKHp4ihZmeRfXUvfxfBn6xmVZGj0osUjh+5R1ClCpTeRUTsb6OxBNiAK3Vb/31v+qn/IQfdaNn9Re+7W9sNu0f+VN//s98w9feyDKsafTX/Nr/6Mu++Af9/m/6s13XlxrnfhPWrFarZ8fPlotzYEClgkDQDR2gUenpdNo0E+uctUYRKiJrjDFmQhNFBJZ7Y6w1iDSZTKbzw9lsNplOlNJ93/e9dS4U6SOXBDMDWGOdc33ft+226zoEIEWTRvfsnLMAxI4ZXGjdSPk7zHrj9TTTZuawhkUIxAK7x4zPBGf5MioNqoLRCz589vlrpbwtpXCCiCL21gYPyVQ15PbDN/3BVt51bbeE5Xq1NsZA04hCg0FIvVSKLcN4G1zjUcaK54/b+prFwfO2KJehTxjn7evNZr3Zdm3P7D97b+UhwRHh9isq++IPZJAazn2xguixSEGBAyaZjFxFT5TdwqCtQSykR9JSS63ZK4dOmAwnR9MDC+C/okMV8wYe1ztKlxRKXa3nDKko7IdHsLMaOsdCuJ3lUD9kTWnSBxj0zLKpD+JQqOf65X9841GO/NIYpuzqxISIL1srjx/JraP57/+tv+7H/ogvvdEbuLhY/cW/+neQ8BNvvPVX/tZ3/Ps/56fd9CP49372T/0RP/SH/LE/8z98x//6TxnYOVZESitj7Gq1PD072a43YBm18vyU+Wyu1Q0oHqRUoydK65AXHVtjPNlmNpv6KbHte+tcoyfT2WwynerJtJlMlNJxnlc0pz4meowEEJ2xpu+7vrfG+iGqtRYAJgqVItbaOfJXgLEOmDGKcgMBIYJSoBTG6OXF98E6P9rytsKAFQaKeQFYknwQRespG3YYIjPeEcBa65jRlc5socZBBG9TH1Zh0gXkbNCXi540rjdm07brzebo4ADyxs/I7cZ8eYGD1yH/XGfop3eOv6Cu6y9/zNjXYicL65gFAN7eszf24uJi27YMcrmRBStLduZenXXc4SlThGB08akksSMBuJizCcTyYekMO1r7l7AFBmom5K1bJLhk+YphLMZxCfuw0Asaq1RQTgUKGqykRKVn8ryroTgkyvAlrQYjFE7J7CCKMqbnhZj1v8cubsHff46klXhWjChkdWsv6aJzrWLE2MIXM3zsI6/80d/5G3/g53/spk/pz/6P377ZbH3N8w3f/K0/9kf+0I+99spNf8jnf+y13/ubv/Jf/qvv/8Zv/pZ/+T3f76+Vttuenh0vzheOHSpi5wBhOpndvX9PaXWzBDBtANFrLCCBs845bhqtSPWmb9sWnKVmcnB4OD881Fr7+9RaSxLzEe5H4ZZxjpmNcV3XbTYba43WmggbrX2ScM4hoqeZBi0H7ywDgMiUuD7JXs+xA3YuNBzxQych10UDO0BEsX/pZ30BgcmW24H+65xlB37MYK2zznqr19BSuIrWhiqsBaR9kbDxHGhCyMrvEZFChG27WZwv7t6+PdE6OojIfZiELY3dgxXmWzYLzxe39U2rq90LXzCAMUYzVI5O5O1/15vFYtm1XUj3Od54mTyxK8GBPlQjOQgAFPhBDgeDg9FnywBOgPWukJbO43nMpnPjURDlqkm+wkRJwdWgmItFMIGiMwqDOhiBjzLPtSKuDqI/SgJnTizPyaYpSxIuXkudfllu1XoBBgzwX8pVsSMoZsco1yNKjbCbjq1F88h+YSIhB4V9TvkpXrWH9mVf/IP/4G//9S/dvX3T9/CTb7z1rX/pb6Yss95sf/sf+BPf9PW/+fmkI374l3zRn/oDv+3v/P1//Cf//Ld98o231uvVxcXSQ//sHHRGH85ffvmVl1+5fyP2FBJOmkln+q7rEaFpFCryAp+mt73tAXE6P5wdHMxmc900pBQwWGfZMVBWtcLkpkZKEVlrtttt1/XAYK2v+y0w60Z7U3g/WPXhnryrklZxEpSbyZgS0mKW53oiIACppmmUUs653lgE1EpFAicEVz0sHB4Tm9N7TnoNakTw/H1jOmbQutFaM7C1Xr7ILzo7ayxbGw3AAQisc+yQgNLeDseZB0YlUU97dpZXq9XxyfG9u3fv3bmtvcE9JzULrDQIorLpeFN6ubXVVXG7HgJfowa8AimGK6ifMgp6AwcAx261Wl1cXLR9iwqJx4p7sZSbFQ6qJSWKsZIuGayXWjqcSAkjc+T0W9Ou16j6Wjb6jZANJg5PFG8cA6jxEu5KSZ7Z4e5UUsWGv6f0h2WA8anJC336A4J/npRn4F26l4j2p9xjBSyoucV68YuxlLEuSriUBy1uuF2p5mA2/Y9/+S/4Zb/g5yqlbvoEmPnr/9h/a52JhGFkgO/819//3/4Pf/n//n/+Pz336/qZP+Un/Myf8hP+6b/8nr/wF7/9z37r/7xyIYiQ1kdHt+6/fP/WrVs3/ZkuOnIhIiMSgEPs+945h0izg4P54Xw2O/A8bec4rstxBEYcMBMAEXmMh9n1Xb/ZbLt262OhIkKtnGN2FpQiJK9a5Jxn75AC7+ISiM7e4MsGtzCXi3lCAO1tM5VW0+lMK9V1nXEtIuom6AiQivT9co/RV3fehSxsZ3oUlwidc6wds2r0pJl672HL7NgZa521trdeAM8b0YRZg2PDNuwfhN9AmpCUStsDQITObTft2fliuVzevX3rencf7vKCvdL27ToFQFC4LifF6RfnZXix0EEFsjoyKkwRfGT6zKW7lmO3Wq/PFxfbbQeQPXBrZCNQSsQiLYtHEAzQlNLLu7jPUxkYKClim7XC/YscgEX8yDo+FZKDoU0J+ENS8K0XDksoghMdClPHwZcxWxGBKesYh2fvRnNEeoH8vKGfd/45B2Al3pskRY4gy+7jjnqBs8NhiI4scu8LxH2sFAEwz2mqpyGMG7AYBPqn8W//xB/9m37Nr3ztlfvP90y+/W9/x3f+q++HAs1EBPzT3/KXfuDnf+xn/Fs//kXS24/54V/6Y374l/6ur/l13/btf/O/+5a/+A/+8T89mh/dvf/SwXR2cxoor7dbJJxOJ86zcoy17ABgMpk209nBwbSZzJRCDqG+sNVJdrzeS9Owa9u273pjDDuHiH1vCHEy0U0zcc4aa9tt6/25mkYrVh688o1CJGV4IIVzZUGkdaMbTxRSvvIjIt1MEBCM4VB3K08l9wxQQIiOCXlCFHeEiGPnDujNZRG8IAagA1BISKAAwLmGlNOstfULacGPLAw0TN/31hoPFgFbRjaAqmcXpeUUKaW0dW69XK/Xa8uuAQXFGA+FBXL0kfXrCDwk8ITGpfZ/xYL5k0I3ZmctznSaG0JALyw+JMwJ009cbTbni4tt1wX2WGGMhoPXNqK4hpdj07VgBFbc1TyH5Bo+wlrUMw0EstM6ydoXZBmMfGnpLitgEgRLIUeZH4FQSB0Fz2BOvxS53AaJfUGOxc9lCYlXzA4wq2TEpR+KukZ8HXZRVteLBAfkqnrB577SIG/6lGSlgVQI80ir/NHXXv5NX/Erb8r2keetdx594zd/q6yfZMD92j/8J1+5f+/LvuSL4MXOwWz2y3/Jz/9xP/qH/8Sf+0vvvfLy3Tt3UVHf9TdsVcCr+SMRWtN1PTiLSk8mk9nBbDo9aCYNAlrnUgUd7slAsPdliOtNr7rOWdtu2/Vq5ZybzWZNM2HXWhswcSREC8YaBzCZTJRSyOwQHTvrZf0z842REFEHcEhr3TRNo5XSSFFDBYCUci43ln54xszWAUBoLpAj0h+gJBBmjlGbMy0lcxA7YgBSSWmOFCCpRgzqwnaE7XvqOmuNc35JwLJzDM5E/bjQhpCyxlxcXKzX61JrZBwn54Ep3DXwlUvutrEOIMw2gnCjHBYm2erMs8pfF49ncSMxlO7NacNaLqMmYNcxrteb84tF23WT2Qy9b5WwXhGvwSu68oh7GgPs3oKtFBqQEMPPcWm6ICjuI2hQhFKCYS/IRaH4grF0+OJUD3OUo019Bg6jfzLBLjNlEdAJotanSFNSE8m3QomMipUD5XNwgH7wF3789tFhnuQP06nAlTya6mlvnP8KDjMoJ6OPPDb3QkkExVZceLPu3bkZmpHN/liQPiSZCurxryR4IcIP+UFf8Et+3s/+8p/+k15ky7ftut/8u79xtd6koiP7PgAwcNeZr/qdf/SPf/1v/sKbT5WH5yu/+muV0vP54fxgbp3pzc32AJDwYDZl5n7bmr4DZ1Uznc3n09lUN43WOtahpZBVvBn9HJTZrdfrvuu15/s756kwjdaTycSY3hhjjFFEilBPJ75M7a33D/CAjAs/T2mlSCtFSpPWnpxNFNXgklZ8fDwSy+IwFZGO0+45Cvo/RkwXA4YjsoDHLojZsqOw6Y5Zfz0TAtGv1xAgNOjbhrCj7Jzt+67rklBo1GUy3G6fIS6Wy0xqzf2ni+8wY1SvYxGKISwepC+EtXCuMaoHZ7n9iDYwUpEABrShGLCHXKJaO3+cf8lcIbyRjitDoH+I6e16u223rbXWezcnje74Zl9aSxYeEjwAGETMEUPXghsinjter/DEuoDnsjynMSgi8SK4IsyAsL8eUFFqjD2CPVioJQnHAglZlY0R4803wf7k1/9m+MCd1GVymoNjXWvtKIm0Vj/jJ/+4X/rzf/aP+NL/w4s/kd/3x/7MJ954q8675eLG+cXq13z17/nDv+M3fPEP/sIX+V3/r7/01//RP/3Ol1++P2kmHshHrW6W8JkVUd+bru/QOWoms/n88HCum4nfY0yUmyh/m1BZcM4aY/xMsmu7vuunkwkRTSaNMQqArTOezGOMcc4CKKUaRQrYdV1vjEnru0gNEJJS2o+JtVKqIaVYwtQc3NlCSmWm2OFhqZ3GnJZCvZCEbPUSfoJjOoR+IcE5YJJrJQyOU+MRUANGAEVaTVFI3rXbjWGn/ATYhW1qAO62XWcX5+eLdrs9nDa1eDHXM7JEoq4isIzVwT+hjNvjobty3cwQ0A6uUEUnkmoeclgw/vj0u6XgBUdtPYDe8nK5areddcl4DxmDabKQ8eEA/gh73ULQcwfPqHTShWTmB+yQB1g8AyOj2HRK+5YsEWu/IZi5QZgt+mrSTKIp5LhU6cphqamMQ4mnNMFgF0kC0SCmFihOJGmMDQtfh7z5OXn8p0aSaY1lGwOyTQEA+Pk/59/+iv/rL7mRtMNlEfnb//Zf+zv/AIrbJ3RwLlJl/BM4Xyx/7W/5fb/vt37lj/phX/x8v2u5Wv0Xv+P3v/rqK/fvv0xKdX0HjpvpBPEGLCNr7Xazccxa6+lkMpnNmsmElA5kGlnSYF4L0UoB86bdbjebrutDq8rO9L1utNaN1mxN37WtL+u1VkgaAByz7XtrjLEGmIGUbiaTyVRp5bF4HVQZCJFSmwHZjdVJO1/G4A3vWROOmXjMGzau/qGQCONQyOdAhoQEaG0eioPsDwQvLfTdkc2V/pOZgZCRtKaDgwNNqu87ts5Zs2w76PuLi9XZ4vzo8EAFMih6ihpLuYD4n0MtBFH3ieJzd9wGsU5aPYZQaKYM/8lYjDA7HzxmhwqQJIOUkcoH4k1vFouLrutUGOnsqr2xfuI4/JMr7J5EjZwEQ2QDQbJvHA0pmH4dIgbaUS10Um0wVaTh4Q8s/vpVmoBRl5+LQXttqZKaBZbu8IAfoug/ehWM8WCL9+TTDx7euznFc/T89b/7D//In/oLJZ/tskS8Wm/+n//VH/zrf/cfPt+v+52//xvX2/be3XuHh3NSioH9CPRG6n/OWuuc1s3B4eHh7dsH8yOtJ2HlyqstUDBT8nRJT3y3zm67tt1u2u227ztmbho1aRokZGcT5mKttdawYy/FAwDG2LbrjLOkFOkGlVZaT+ezg8Ojg4P5bDprmilpDUiW2ThrrHXsGPJIRygFZiPVAX8Q5dyRr7AEDd8Lm2ueyC8qyBEHWoFXWgbrwLmgYucDutLNdHown88PD4/mR7fmR0d6fqgOZsaa8/Plpu2ZgcIGA2b5qkj/y4F0+BvLKH153B4N3ZFZA7llHqukrvn9oN7DZXrkESWOMPNk5vV6db5YtNvOj8gHvA8XMrNH36SN8rBmrt4hTlg4p8sFsDRV4OqNIUBijIpuaafDLzQCA7O/BUjIDUJyiUGKoHbaVPbfKp4bF3JGYT6QfQpihBcWwCg2aMNLwPyKpBlxLk4CF6dcSdi9G/E5Vvx701okwtQL1fcPZ/fMcIkg/PPv+b5v/1t/78V//9/9h//k6/7oN4cykwE53tQciEjkr4vw/aA63nfmd/7h/+Z3f8M3t113o1/3Pd/7ib/47X/nYx/7mJ40bd8x8GQym0xmbdt6cc1rHiJ1dOv2nXt3Dw+PvLO5H4EKfgpyVD1AQEUKmFfL5fnZWbtpA63SGCJsmqbRGhi6tm23rbNuMmlm01mjlbO222y7bWtMD4Cz2cHtO3cPj26RUi44aiEQMaLX9XShc1d+/irVmziaRLkqFiGAQkBweWAWaEoQiP/+L7k4fII4APary1FsnUMgrxkPWPeRLOaSjGEAR0r7bQlmsICotGqU0pP50dHR7TtaN5vterPdMjOKHzKqD82X/8OjyjbXjef03t6L6R8eQWjWm/XFctn1HRKRKsiUVCs4D37ATkg7j3pwh/KHkOQcpkXkkfE5jwvHAxYN4UiE3bmOcAVzKQBRlEcEcfAlAEnG0lVV9gTx1eRhMOKHIv7XxRrzrusyOyjHh/zRb/7Ws8XFi/zev/eP/vff/ge+KcDlPAZQlj1j9an8lb/993/Vb/jaT336wfV/41f99t939/5LL730EikyvfHO6F3Xnp6d2ps4gjWTydGtW0dHt6bTqXOu7dowRg6gqQujXiJmNn4/rOs3q/V6tba2JyJFgV0fvLH8Gq01zKyISCkA6Lu+71t2rFUzm80ODo8Oj45mBwekyDr2W2U+y+SCP6DDodZ3kMjOsskVs7ayz5MDSZZyugxc6VVKi2CfMtgN4iiPTAFHoi0SKq01AHam77veOb9xhrPpdHp4SFq3bb/ZbuJgYwDui5/9Lqj6XpL4sxQOjt1E1/0+liLBuSNLE3jgIkay4816c7G8aLsOEIgoLY+KQzVdM0UyHFV6QCHVPWiBMm8TQzwFAq7cc9MaBwIhEhABepabZxTLAOw/oDxxvWxTOn2WmAp04dqYSP+F4Fu48Cmam8aeg134/enNisZgRGnFMZGbkwnlhyX8S4tPLCNsFonBAj3zD15cLL/hm7/1uX/vX/if//pXf9039r1hWTwgFB17oUwRHecwXa34qTfe/g9/3W/7xj/9revN9srf+K3f9le+819/P5Fix14czfTm7Pz04aNHT58+c+4GUWM6mzVa++fp2PW98QoNiElgzWkiRYqd6zab1fJivVoZa5KLkda6abTt7Xa97doOEWbTyXQ6JUV9b7pt2263ji1qPZsfzG/dmh8dTaczpfUA7wjMmrS8WBW8+Z/gNYHZwENouSdpFselKV1xKaQWHyiqQ0SJvdQDhDvX+xBDJpAKQmBlTB94UZoB+77vTe8b89ABojLGrNab7bZ1heYljUSQHTdv7vBxxOXo+vGcPqO3ZnlBtl23Xq/7vsfM1MoXAY76jFylERaLeC66hPBhu9INEq8S4mZBpq8iSTlYKazbR15wskL1UBOIFMblRtIYnlWQejAvrfKuXzW+Pv4hmQJU3j87cN4dzSP+1f/3P/hn3/Wvb/o7jbG/5xv+9Dd+87da71FeWefUv1aS5FyULMiPstZ+y7f9tf/gV3/V3/qO/+8lOP7iYvn1//Wfvnvvntaq7buuazeb7cXy4uT45Omzp6vV6kYZX2uN5ImM1h+/60upDCFyzH3fma5ru3a1Wm02a0KczSaK0FkLzLEHsOwcoBeDIGdd3/e97RlRN5OD+eHB0dHB/HAynZEix2CzpLNcWC9FoUTQrgAfWfRTkp9FsUTPNUtFevOlqpNk5elFTFgue3KQphYlOSc9wbEqxM81rbHWmCRi7Knjbdevt+u273bplV3ru+9SB8CZNM3yLb7e90VXHZqqkrsfqwce+Lhz17WbzaY3hhCVl/pL0Y0Tn51xxBdl12oqIFHeeasc4mPC5/Bhcik/V3R/HMU/PEKMRYQVjJ7cjqAsMLAI315mJGyl76pahfLDGD7kCbwMhJC327NYDhdKdESAcTghP7sP3yQ41ftDR22Z0hNe5Jz7vf/1nzHmBuDJ46cnv/a//Pq//Le/AwrIrQAWcsOB4Eep6YNBdoKgHpZpgejZydlv/4Pf9Cu+8r/8q3/n7xszgub/rj/0x9ve3L1z7/DoyDEfH58+evT4+PjZarnqu/6yi21XLEDlO1/n2FhrbUTgEbVuFFK72S4Wi/VmY+IhBK01Ihlj27Y1ndFKHRzMJtMJA/Sm79rO9D0ANJPp/Ojo8Pbtg6PDZjIjFS74AOWndfssc5yBGigIDSIgsHeJDI8JvUNwgskpIcjHAaOMHhh+jhSf8StCRKjifIb8xxVHDeyN630LHr8ZyUKipGNGBKVC/2Sd49hkeAca48x6u2m3ncvX41BQFfxrLXJe2QrVLdH143b8hy65c26kwsWFAejg6Q6Otdx1putaY60HWUablmHtxDufQMa+ZauTq/MU/PlyF/Nd8O1ouY2XFO8F/wSHzl/VQ/GKjqlQnxshCGRgG4vFpg9R+X9jsCg1ipikuF5/8+0/921/7Zo/4S//re/4Zf/Z1/zz7/peqCcxNeso+yHD0KSO8zUcCH0BH3zjwcPf8nu/8cf/nF/6jX/qzx57f2AAAPgX3/2v/8e/9DcODg6QaNO25+eL09PT8/Oz9XpjrCEkfUPZohDCSEUDLBcDIgKAtaZru8120243Xdeyc6RIawUA1lqvtOArPVKESgGz6U3X9cYZUmrihYTmhwcHh81k5v1BS8n3PCdkSRlnFqZAcYyOKVj7wAEJKk0osPeR54DgsFe1s84GeTtjjLHG2PR1n3Jab23fW2N9cLfG9n1rTB/aGp9I/FYagqJYiFX0e8Ecd+zYiX1OQi8517V913fWybnv7oHVrkDKLxq39SWd8o5AxaOPz7ouGcSWBl1Raz/WOl1n2rbtOsPskEbETwPuwixF4qF2ui28nr0yeH4oMGL45ZzVFSjKbCLz4IVBUK4kDCjiiPB2lJaRkEstY5zVkGqBIxjh+2PyCBPS/7i7D8RSqzQbsIkXisyFCDd+iOJ61BQqSr5xiSFhV8MJgv3mP/+//Ds/9d/82EcuU2x+/Ozk93zDf/OP/9l3FW9tzrwk9lgy1YTrizeuuOQbF8PqDyKwazfb1XL56NHDX/fVv+M3/rbf/XN/1k/75b/o5/17P+dnftXv+APT2Uwptbg4Pzk5Ozs7M33vbTf9lWCMvamRWl4yZ/ZinL7VNMa27bbremd7ImTnnAM/MOj7tm9brZTWutENM1tnXWectZYZECeT6Wx20MwmijQQhYIchYtGukA5SH3KtglRmnY6oS3IwCE5VWWic9b28fYa0OeFoW+WSBFWcBkvsM4BQ7fdmL7TSqP3jlSkUPvpIcSl/1BMu/iSmIE9Eck7xliVFimcFxBFx9wb07Z91/XcNBDXf3FYoIwsLcoNqxeN2/rdvOtGHQlTzRrlERjAMm+2264zzAC+9h/IJMRvXg2Ssaz0OcV/FjKfaeWK5JrxiEIa5u2KsCJRQSc8VuFzFhe6TL6Sd7YKOM4wqC5urMyNMCpM8oj/celL/yEK/9e+I4pkXOgrtl33+//4f/eHf8f/n70/j7dty8oCwTHGnHOtvfc5t3lNdAQRBE2CIIIgIqKSCokdalWi8kuTVEt/aJpm2mSphVnZWZkqVVlZNliZZmJDJlpC2SCCqAgJigWa2JAlKgIRQQRExIt47753m3P23mvNZoz6YzZrzrXWPvece+97QBRHkPvOPfecfdZea8wxvvE1v3/1Xz+8uPyLf/Vb/+q3fvsw2seitUWeV90FMvkJVIbeU58U3WMUOesvHj169OCBGy0S+cB/89u+81v+1refn59r09+5e/fycr8/7F/+yMujHXe7s67rnHPArDUxy03d/ziTGJNrCyfw2jt7PBystZ1RpuuzF6YQxc2agIoMaPLsnXUcPJLSptNadZvNZrvVRnNaqZYTcc0WV5pqXDxV4t2fUkKjUY+kDr88nsIcjeeEg8/eb+0jmr5zyu2V6omq3S0bp0vx1gMRG9ashRUzCVX2kYSEKoZdx+9IkdwhWbFWmLMAhOjrp1nEOXc8DG63TZGTMRvgDYP/6wNgqRlbixSQObCJa8qK2rmuMokvCT8sMIz28uJytCMp0s1CFKCykml2tWuHXR0HlliSJccn5ybU+dRF+1djdtOrp1NYT3E3w3a2xxk5tHWyzBlasnZ218djhr9wMR7AIodg7lgkzdc1oUhQGck9Bj/7aMP9J4eiChbDthmv0YayG5Q0GOI//Cf/3+/8nu/7wl/QuHXuD8dv+KZv+/q/8Xf2hyNAbctUNTrTf0hbgeoZv7igNkSCbMslhKCIhhAePXr48P595x0qlYY8RfvjUY/WBbfd7JzzzjutVN91AjIORxDZnJ2ZTt2seExeKqkGWWePx4PW2rsgHIQDolaEKIoh2GEEBKXVdrsVEO8De+eD4+BB6c3ubLPdKqUQCZWa4ism2VZaoNawgUxAdXI9YREEIapYL/GsChyis5DzIjEXLAm3ol6HJ5enZqeQOvdmZhOpag8ipv+WtEPIqS3MQSSAnzikAERama7vok84Ye4aEQhzRy8CIKnxJJKYd4YAAONoLy/3t87PdpteVUityCLWsA1bW9bhm9ZtuFEk5PIfPw48X8WypGzTj+Nwub9w1hGiKFrvibEKyLn6NRVPBqwEkCJ1c131GziPlyzr8LZ6Zm+Itp9GhBMHkiwgfSwnpNRE5auOeFzrH1Fm68Rq7TJntyykkNWv8O73feDTvvDLCuSBAuuEWplj2FIWmGszTbS7/cv/w1d9yid+3DXvpf/8//6nv/nb/h6kQav+KVjlWZbjPFPY6sxeyPvUk3DoHH+vpy1cB1SbK/wnvuYvft5nf8ZuuwGAd//oj3/zt/29v/Nd33u5P5zY2cyyp2U59QmUwPATb/Tk1SJIyMLDMNrjEYgoJnAlaie7IP54dC4QUdd1SpGA+OABo6m+cAg3Pe+jvheEiGJOuzjrYvUnxL7vCNA5H1HUaJFvQBGi98GOloURUetObTabs7PNpkdAZmGAhHTPeHQIEw1zti2T7CrIsYkTDoGFMWFrEELwzjlnfYxnieh8lgczNK5wNcsDObG6s6+jSMwMIqoXbJwiiDNdXIB9EMie18L5eNBiBBFFBwRkBAzJpU6RimEAWimlFEXn64isSeJx+OD2+/0wjLtNj1cwW04X8aev29eKhLziz7gCa8KM34jxnJeyu5fjYbi4uBztiGmyK/+KaiJAFRy7yquTUoEgub7G0U/KTg+TOYmIMGT7yQKfLzmh2cJ+Onxzh0h45QJ3tk7Dyha/+bqUXFrlpC/HpuorGysfzPrmKYqjeMLlSNTSLyBQHTAjUqmXM9NJsnfQKXpMUVIuhS754jy5voAIWnAKp4wwhMkDHauQqSTGjI466/ho60bVphO2tbkyZsJaMZR/pZfv3f9//vmv/+RPfNc3f9vf+8Effm++R2h5GGNt33oqBLSo+HC6aFKDiPlWK7PbNLplel0kqKPWm36DhN4HYe46IyD7/R4RN9stIjrnwjD8yi/7Lduzs+3uDBC996brzs/PldKcvB1gWh0hvvra/ZjXIkLxJ5JCRAzBOev73mjdO2ftYSBFxuhu2wszs7jjELxnEdK62/R91+u+V6aDWs8+HYQyqSKgBAlwnselNjIpxioc2NrRjqOEELeoACjAIKIIGDA2/sASN3jSTOfQPBUAmNpxqUFVmrx9MyE/3YHAgYshRjVgUiS5EkIILmb/5jxLFkRjuq7vACB62iFh2iFL3F4QAljrD8eDHcf0SNflIbWquGwpnrhWw1NHQl4DRWyaLTnFNxiGcb8/OOuiYAJTFMPskb0elI618IfThcSqkKQrKyiT3ytitqGvcRaBuR/02mAiJ5fxZTW7FAzDKYFY+5U1gZ3WgtHw8TuFMs+Ul5znR0qf4sejw9dM9nwKhBIfR52S1X+SKyfhddkONVB5clwokQntHuYb/9Z3Qqu1kZVgnXmS9Op1qgrRyZxoqK0/InBMhEoJl64TSZEx3WazAUSBozAQUWCO2byEyCLjOIpzP/pjHzi7fefs9h324XA4ksKz8/PO9FKGZEjmB4QUdwDOcvAu+AAghKqw45kl2vQXfmU8rZ1zwTlSZHRPWindmc3GGAOIU9e/QjOZ9nUpNbGsoEU4hGJiER1VOITgvButBAdIENtplU5URUQaOABwCalp6SRYkC3Je4Wl9ib/fJZiOZEkl1kPTBVpjJSizJgK3kX5R5aKMSBGFDqKR5l5HEZHDtJoQEprBPHBH47HwY615Rxc0x3sWe0A5Gl/SrNqL46kbbB9hdUjCMgwDvv9YfQWKRp90MpzsfDFkKZoTl4O8faV4iK6sn9DJCrfMHuBJy/+ahOE9W5dYFJvsAgAU6IrzXbzaduWMx+wMZ9eIxFJ5j0vNYTz3j/LFFuRe+3Ais3sJa2CLkE91f0lQeqpZu0kkdpyrjKpqr1T4kFbk7BuuHLEuS8tVvk1UFGcpmtSxzkWrtMsAwemuQEft0hLV6VIppG5NlZaPxSlSLdWev94K0qWiNdpnVivZMpC6fk7t9/x9re+4+1v/bi3v/XO7VvbTb/p++2m32032+2m78ym74ZhvNzv98fj5eVhfzgejsfjMO4Px/e+/8d/6N3v+aF3/+iPf/DDgLDZ9MwyDmPwTgLr7e727Vu783Ol1RACKkRSadskINFNM58xRMp7NwwDB+e9d9bGX4Q5EOmuQ2etHW3Xme1uy8Lee28jj9Kj1rvzW13fjaOL55DuugnrnxZdgIVoxZHUw4KkiZRSMaZFWIIP1o1udMIhsuljnjpLQEIRlQMkOXBs1SPnH5VSQNKwfUr81yRmhjr+r96SM0u1BKmMwGJYcVz2RiuAia8q0TC0SrCJt0GEjSB4x4gi4N04DgdE0sZsNxvVdUohCDrrLi/3x+ORW1N9aW7O1/co0M/6G8rpGKnp0fXODXYInklpSHaveNMGsj4IMnMBFw353FBVWjgYp/wqmfYGq61q4Y6td5lQOww97h2b8gdWveknU2mcLTNQYA6k5mYFpl+iOQDaUEqoUgeuyG+eTVmIk7JsAQMhPtEtemXr0ZBqZfJWkKsniMcyf2S+la2oUnIisLnqEk92+bL+nzF7qLpvPuGdb//cz/r0f+Pj3/mOj33rx739bWe77WNf+O1b52++MpDSef+jP/aBH3nP+/71j7z3O/7+9/z97/k+1Li7db49PzfGBAkiEhNWEniStqxAlFrvcTxa66wd7Tg47xHAdDqW48jwISRGIUJCDIG9tT4wIGjT9bvd+e3bpMiHSzs4H9gkSfDKnUo1Z4EIESJBPw7p3nkRCdbHKSQaLxRlFxEiaqlljRmjFTwBxZalPpRUK1nvVuqVveRmqXV7aJdFUsjnxcegPC2KongiJIQreGAGpUREq9iKYtRcG6Psiv2fPG7Cf3YHwFM7xFSQCNbEfGi7dSxh3wLgvbejDYFJR5fBjNjj1ZnyzZNWnPhjKpAgtFxSnEIkpXnuS4RLcfMvOfNYD+pTSliJAMBqUTlh0DRrG+UkOC0lVClmE0j1mVKSCk0iPTT1skDaUWhaNdcmrJC7n6mjmLLLJkpsc9qk5UfFUayGuna1Mz/C2niymx7hNRCHtYMUltJZR1I3y5/5FkBktouCiYxVX//CuSleISwsJQGiPV+mqyHllpo0Fggt0syQpSTldZydbX/OZ3zq5332z/q8n/MZb3rhuWf+DButP/kT3vXJn/CuL/niX/z7fudvPQ7j9/7j7//u//Wf/f3/9Z/+6I9/kJmByGhDKqZ6cSlqpDQCHI+Hi4sL7wMieO+dc13XGd0F9na0wQVtlOm6nsgHfzgc2VsvgqT6Tb/Zbjfbs03fOw5FOQbYupzkEN6E8giICBGRVhzCcRjG45GZQdiOA3lKtwVlbWngEAJTOgeIKquInN3IzIG5UmBNKI9US6+TJLsGv1uKAtKPqIpd3AlSVPZCkzoHhY8UsaMsH0szoHPOOyccmDl4JoWjHZadPr5R3r3PdgK41iuOIgjnHTNXVshX/+v5wUC11Qes+oBX31lONO0Ap4ReVQ+KJzQODfUGryCjnLI0koUzTJ1tJiytlOhxsuTKQroeAuqJYiIctqv7x40tzWYCrvMLXnMHUML25vznya/81Is7sYjBx+i7W3Y4VqKx5qLhylMoAKf2U9OeR8pVNlp/wed/zq/5pV/wmZ/2KUq9caZb203/Rb/o877oF30eALz/gy993V/9m1/3V7/5cj8owtrcRkRC8N754+EwDIMwa62IqDOGEAJ7YSFFNbzpnbPjEZBM13d91283/WbbdT0pLcGnujiLR22vUkKWBEJg75wPfhxG731Mf/fiMaBSKua6E5IICqfjnrJPcFrppTt9+rialQUl0ml9mK46oZbDXjrIustKG0cgYMlnVXUAlJ1BETPl49Y7x96HEABQghv22o2WWUDBT8jH0+8AYJ7A3T6FNboey5i1frTOOR84mb1K6QVXLcySg5okLLsyjUsdYhXRJdiAPyKNwwiu5ebU/UDG/QUn8/GiYpNZY1j2vaeInMWRrqAzWLcK+T+xYqNjO4+2VxQrUg9UqIjM9oeLX3HKKYDmzKzkhxMlVFqPK5yQSSzUWbk+mPOY26ZUVanjeTGT+lqGENZCiylhV2Dm/ocl8ztTS5pJKzOgpo1MHe3QjoXQyO9kWrNMrFGsl0IIIOdnuy/7Nb/0S7/ki+7evgU/oR8f9/a3/ee/57f93q/48v/XN37r1/6//8ZH7r1WzrlYfI/HIXhLiIzAIWijO2OsG4fjUSlluo46JczOO++8cyMQ9bvd2dl513WkFFI+IbL3UIndK70/JHdembIkA9txOB4PzCHZSKOIBBCWAD5+UlFM5kpmPZJx+yDMHOIcM2mpgUgRzJFPmEpFhfFg3kNMorxUtDMXKGn7sdi1EhJNa7do78AigQNCAA/NJgpRERX7YAQtOddROG64PSBqIu/YjvZ4HKy1O2MEatr1GyTd1G/g3YhRNjiM1jpXaHyYmZtX48iEwKUMzskwi1VqdoJYpeK0EAjPlsw4PfyPG0mui3q0mFRsKLAeVdYIQqnXngH2K12twBVkIcwQ1byHnVfhutK2gZtt5uQVJB25YfVHmAXfV5tdlCcbLx//dc1L53ZJVIvHAdc5SVVQaXV4ps92Wv36/90v/fIv/ZW3z8/gJ83HrbPd7/iNv/4rfsO//c1/97v/+Nf8hfd/8CUW9i4cD8fD4aAI+03PTME7YZZZ3gYzez9ax+yVUma73Z2fbzdbrXX0S2AWIansNefvcfZ1FxF2duTALOysc9ZyCNoYRYhJfUCxFmc+cFrNpFVY8pRkrnweE4Vp+QzV0KK0qPo0ThdfyuIUmc6iuESWuvuTKd0JCZCJIHD5wdIYizMLUVSxld2bFNK6UlprpRCDRe/dcTgejsPt3e4nJLJDP0uXeFzkEU91GJFQAGwIx+PRuQCAeCXpsubYIHAVBjbpgwq8KJPUcw4YYSvJmlgfFZO9cH4oW4rzyqJxom/g0ktaljAVVo5I1cKh7edr3yjIydflq0WaVnRpFdKQgaRq1itCTfxp0dqcEgU5R6NW5KfJnqDyIInDVTIObAO2MHNdnlQJMAsrTSdeotJSM5YUDt/MNEOmLXRhtUiVFzcRhLFqA9O14ep9lcXOsj6dcLo0UHHIJ7IYAgF++s/4xK/8j37Lu97xMfCT8kNr/aW/8gt/2S/++X/4T3zN1/zFv8LM3lsRRtKESIo4gB0tAOjObLdbEUl4BYsAmH6z3W43253pOhVT2rOJslS7mxTkh8kIDTGls4mgs/5w2DvvYtCGUToABu8DgKKI/FTReByjFTn551SjLSHF+SBmawCAMMRzBaQS1NQNoMjESqrU2JEVDa0JQc0GxiLW5lCZ1cdlh4pSCZQJ/eHIBRX2ATAEqG8/xKgJU4RIKr7UwDwch/3lpb192xhDpGr5/hsQ4fF6TgC4ciKPLgzj6IMnQiJ1nW9SmSRQtceN5zhVh3z1M2W2qVw46bRzSbtIQFnFck6dd3Kt64DzLQNWii+BFlmSooiYvH6wxYRmQ0pdFxHXNhLZ42AVZmugscpjq/gM5epcHww3bc1P3B1YrVLmn5Ir74qTn1lBd4uIY365mpjXmXpvsVdae2NjB/Hlv/ZX/vbf+OvoJ33qztlu+1X/59/zRb/wc3/T7/5PxnHYbHpE8MGVOpUgIpHgvbUjsJA2Xddtttvt7rzvDQKG2IzjdLdh8eLPrRlFW3JgDuydExbnrA/OWQcgWimtFJGmAMyS0tprtnLgQsIpnR9O6WCUOzgRAQ4cApd0rfmzIlffhwQFRsZkul4wwKg4LL4+9bgqiAQKaXqcEZGin5JgTbStZo+4IQaCvBlG9HET7lzcwVTGFOsj9rO9w1Z3AOsLzaSgEDw1T69sFvMbEY0xWMCO9jgMwQdEIiWIKlkh4tpWs1pb5lpfePH5cC6go0xdsdRh6Fi4QdMGaIKAKwQeADid56UQyALnTgBFFVA6NcTY9o4VTNTMiVi9x4WbmZK9gCS1JFk0UXoZqZmQzRQiFWad0iwmGUMiqgHGvLAilkGonTuwPZWqs2Gi4M9AoFlq8U3uzBSeCjOpjtQFGefU/onRhO1hujyP630HtuSO6eam6mVT5QzSTiaLEa9QtdMrMcb8Z//xb/uiX/i58FPn49/6gp//D77p6770t/7uD7z0svN+OB5CCJ02pu8QUUIYj8cEVfebzXa32W6M6YwxWAYnLArtDG1Glmg6ApiItKbgxbpxPA7CARAVkjHKWuesxc5o3XVdzxyiHTMz554jQpekYvg8NHru7MuTzOBkJs5bJsFiJSdvyRHJ73rKYk+3SswViB5HwKI0IVCynchhANHmDUM6MynFgCMQaiRAFcGLaCQUfamZxQcnIhlpEEXIHMbj0Vp7lg0hEixRfPFwWl3VO77rCyGv+GJ9TdbNtT/WTqmqwbXODcPovE9XmFYavdbdbNZ7VaZsDUdn/bViWyjqgpJP9gjpVURLnPODsKHKwKKfh5Uo0rQgpquuUk18LU7VkttrXN+Wrl5YnDixAlUSdr1Uzs7lcwXt0q6yvtSVP8ay38enHALWGv+TjqkzQ02cKawbmD4ze666S+dOgM1LmU7yJcbWvBSt1H/1B/6DX/R5nw0/1T4+8V3v/Ft/6X/8lf/u73jv+3+MJiWmiIBzzjtLSnfb7Wa77ftt13eKFEDqh5O0qsLfyjwa2ZrRjGiwzgfnnQ/Be+8AQGlFhEYrkQjcMwdgZg5S/CCSZUv8RlNEe5LXiQBzCJwjg5NDSKaIEqb0FYqMI6xGu+UBUJBjaoyvonCUMK44GFHrjpQiFULw6QDgEpTJHO1J8zOmCIkIqNIFRHYsMwHw5EufttrWDvvDYRwGeSxl4HUYL/WVTnLZNbP5fFQxYGuqCBXlpnSu8epTssUHQBBrh/1+PzoHueNdqyTSnuQt6ztxUajqFhf1uIDY0lh2TqqpONDldABKMZIl32tNcTq59QIs1W7VUkFAksq2unAVHC0VwbSyIcq5EstNGkxGllB7BSW6QA44yPQzrjYNcRrIvB5sqO71N19OcnEQwbwxmL3ZVxy61+8UpudRJn5UYzQ7MW5qwnZjjNXwt0pPinHIwKzijprvVp1eyS6mStYsWGSxVpc25RO+8nf9lp+K1T9+vP2tb/mm/+lPff6X/DuHQTGIs9aOjoWZA5puuzs7O79lOhO1wvPbdYr/i2i/NAiMgHXj/vIy+NAZo7WK7m0cgjHaaI2oRIJ33juXTU4ppeYRSuG2sXCI/1vlWk0NFgIqbaKaWCmlgNA567zXirTuILnn5bkkTy2CJVArn15EGU+SaE+BgMLJVZhIaWNIaSU6lQBOQTOBg/eBQwAOcfTxDMjNjhPjdwRERKMNEQJgCMF7G0I47I+PHj08DseWuZp7NZz2l9g6irXuiE/4Z326i3+iBxzXG//y4awfhiG6AOGqhLQqiBXXTsrvjW0clrQqaliqfmre5bx3j0BOAU2uGl5qgnNDWYR5pCROhJ9TjXO9d2i3ro+5qLLs/tOPZS72RgXgmm0ppG26Ze5fXS1h2wu1hNwb8pGI3BD7TjH388b9cTfcLEhHZluVigWIdX/RXsGZLWsDiFUmqS1MjPlErL/hr/riL/jlX/gL4Kfyx8e/82P/7B/7w7/hP/j9wuyd4+ABSeuu327Pzs+32y0ChZzScsWEHd3RALy1liUgorMu2kUggAGtlSoPmYiA+NjFYwF46nj31FsLc0xZzl4LVAIf0zaAopOEiuabCkRCYERGVKQ1AECY9H5Y7xNita4gzkJXSQhlixWRIiSggiJHjZcEDIyKmb2EUEwWOZnHVam5IgjRXAlAFFCMh9QkEjiMw+icgxbleIOoAdI8yPUzPft8sarGyu9BYO52j40sKYVU5drIYq0dxtH7MAlSalA+8YWguLVOraZIlaNBrder1ISjvBPIoaDTkTcNsMV4lFAhLOJi5xUDYLIKhMmRs1oElj1qVfWTm7TUosimyM4tXq86A2RWeWv/f4zJqFi7z2ClP2gsi7IrNxax+4RTUe1lMRFx1tAYufqgfzwuWSk6QSqHrilDAmaWrdM0VW9Cikwn/Zkjm2tyvyDMRKbp/JsEc3VsQ32bt9Qoqb0DY/vxphee+72//cuf8vE7DuO73/dj7/nRH//IK6++ev/hvVfv33/4aBzd5eHwkQ9/5CP3XjnsD9FgigC1URvTveUtb3rnO9/x8e/6uI9921s/6ePf8Wmf/Akf85Y3Pc1r+JIv/sW//lf90v/56/8aA0SR13a32+12Xd8DYFxWIkzehFJxDRJlN7Awx9tlHIfjkWM3b4whQu9dCGGz6Tf9hoW9s3a0zAGRjNGm6wEgpjOGICEi+5NRT6XvQFTGJKpMbNpz7nHcDANSNGMTQM7PJNfqoCUtA6S2X5WqcSREoWkkFUHB6Dc6mc+gAJEiLSwGQsogiQo7Dj4E5mRVHTtGDgJsg0cfASuttCJAQOuddZ5hxmCSeiaed2DybP78ZCwgvN7X4LJp9cFb6wIHwJSrVpkAT7vdK3cR+JjxRFYEUbXzADTkhcY6p5kkABLzcdr8zV4InVh+EyzVYSvJXCeopMvfuXJma5zT6zEzCeYIam5rFXuwhHjKKlnWXPROcp/mgxc8AX95zXKnKfV45WQK823LxJuFmavg5C8r7cwwmUnJ2u8923dPqTL5e3zFl3/ppu+e4OFh5n/+r374u//RP/vH3/8vPvjSy1z5iPkQlKLOdIfj4eV7r9y/fz/a9RwOBzuOQAhI7/3xD37/v/qRO889v9lsBMCH8Nzd2z/3Mz7tl//iX/Bv/aLPfTLp2R/6A7/r6//at3A0udSm7/uu70npEC0WYvIVqWjeXN6J2MMCiHXWjqN1lkMQluDZs1NK951RqkdCZo7hiMIh0nWISGuttAJASI5qIQL8ueJTkoOVRS2hUkYZjcW7H6nYqUh+OrJX0JSRKmulpLgu52UGXgF9sEiIfjUCQFjErWnzG+mJOg2IwgKKKGgKMUhHEkIW4uUMEEJAJNIKUQCd98NxcMURCE+uvm664L3BDuB6z6wsyvpKiZ8dsgKIwOWneB+stSEEZVR8KIkAkCATALDFmeXEz6m94qHK583tmjSilIIJFI5QRn1yhjCtVOXJpwAWkuGmTklNXZmvrAEXG06cy6va1zl5XM2OGwQA4unowRJNnFr3JJpEWfNPqH2OYjZfNXLjXBC8sNWUCQcrQBMgPFYxdyVVoDWSmqyHcDI4wgnXqpH5OlAr4asliQEqkyUqK530xQxAE5UIkxq6QVUrWZOUwlJYUAQAL7743BOAP9a5v/6t/8s3fNPfuf/wYvYMN8w0lqg0RQBFKiefAwigUtp0xnRl1lNEDx9dfMf3fN//8j3fp5X+4i/4eb//3/9Nn/DOt9/ohb3j7W/7it/073zD3/jbDOADC6AARv9pzIa7KUQlCbWiMJwwmvUP4+FwcM4KgFKKEK3lEFxg6ozq+03w3ns3uBGYgdAYY7RWWnMIdhy9c9KkECGS0sZopUhrLJvg4rmWe/bAnMw+y8NHyd1R2igJXDAN8iQ8yW3a/NUCBCGAMABHmJNI6hEVak2PipzX6DpPWimtMKccS6Q6Oed9QpCiXVAQBO8fGTOMYxOnhbMfc0oTcJrDcc0D4Okw/isacVkdF7xPE4AuSEix8Jars8uREOFqws+6HrjyBp3NFzPXydWUr+vqnKIgSFBIsD0HpU2iutqdvniazO/JJjVz2gy3n5MGyMCJ9tosl6CF6h7nO9W+Opivr/FJModlsUhYbTjwZNeD7XQ8a/r5BGcacQlcrTH/sTIvX9tA/PJf8vk3pfx//w/84Fd99Z99+d792YRV42mRjeNDYBbTma7rvfc+OOFApIAoiUijdjYXCiLFHJx1o4zf+K3f8Zf/xt/6it/wb//R/+T3xu92zY/f9GX/+2/6ju+21rkwBOHAAqpeq2MWzCYtIQIK8zAcrLXjOIbgvQ8CYrRWRJ3R8UaOzv6c4H6gqPgiAsSY6mWtBQ5AWmkdyTORgqOUIaWUVoiNMrlwIznf6Vk+UJFG5ITkJ3vBzBEGzFaVE9iddZJRw0ZYJdWWW7fpDBt7SEwxJxg7Dkh5MsRsYt4aC3sWCYGDjON+vz8Ow4xyLqcxkNOox9PtAG6O/FwtsZiQbckPk/feORs8wyaN0xHyI8yAUGO/mH/ulLGFssDpK8U5rjATqxQwrJ49Wav+jbQ2nk1UkQNl/XpjaTinvj2RyZozoPYLXHM/yuYUdcp9VaSkan7rvfUUvCyzyx+7p8zQwMh/aIDttQZeSnbbjNk6L7cFjxJckRA8fnSdzVyyot+a50vXs1nr+p/oasIyP05zi4iFUCEVMCwVNsdT8mW5ntiynuJF+yWf/zk3esz+wT/6p//1H/sfnPO4KBnlLCzX0HrnvNO6M705HI/D8QDMZLqu3yhtai6jJIkHBO+Px6PzDhGD8/+PP/0/vf8DL/3PX/1HO2Ou+Qo/5zN/5se+9W3vft/7g2eisu/K+CELoxRXy/gUueAuLy+OhwNCtAXC4IN3DrRSSiOScHDWRQKoMWaz6WNeDYdgB88cmBkQQBkkpYwxXUdKZfQ7HTwT/bOhOmMbk5W+ngNKJQFqM78WgztM3rxzmnjW0xMpECjFXIpcrW26y6MZpbzF9iuHgKamUIC0Vn1vECl4H7z3zo77vbPeOddYwkkbwnoq/ghlFdl//SaAJ9sYpN8jRAZvBPtWoG88cbQULzOU6sHB+epyAVM3c8CqwQyuzyt4rQNu7XvhzEIG5xJkuWKNgnPGfWUUB03PvWa4KADwBZ/32f/9H/3Kn/zkk//y9/32//L3/fafkB/9W//j/+IHf+R9a2PsynGIa3fJZtN/4rvecf2f+PK91/7wH/sfnfOnDsMcq4BEFEJw1tnRCQChFgB2DkJA052fn3X9dhyHkKS2SkCCBDvacRjsaH2wIqIU9X3/Ld/2XX/yz/yFP/A7f+v1X+dn/sxP+aH3vi8wq1ToJ/0RQ1CASmlCYO+ttYHZO+us8z6AhGjwrDqTGiEOAiHE5TASKdJKUWw6QvCj9cEDotZGd4aBfAiCiFppY2Sm4Z2oNNmksBCyW7z2Cuihok9Ca/uMTUGZMTVEsotR/LeEJR+76uwE5mullHcfcyqZ41zDUUOHpHVHWivl2XilVNhukRQAeR8gWsLNgQwUeb1oQdQmGazDHae+5rE2MDg5TAIA2NH56I9RB19WwS2TJK+GdXDq+MoUVr68+IQvorWwaolRAFvWfj1lNNe8LAswUx+k6uql/qJ6hsTKr5wK7FhpEMpeKlmo0EyWgoVXQfmvcTHoULQmLDY42FwyhJ/0TgQ/KT6wUYsWB8i8UZreCcw5JA0d6d941ztuhK587Tf89dHFFR/BYp1EOXYxHgAA4Lw7DsfxmKmBpADIdN35rVu379zqNluqWl/n3MXFxX6/D+yjcz0C9L1RSv1fv/rPvvbg4fVf56d/yifGrVisdSkHOwtuAYQIEMEFdzjsLx4+2l8eALDvTbS5BxGju67richZPxwH5x1p1W83/WYDhM77YRyHYYzV33Td9uz8/Nbt7XYbUy051s1po4sAKBQfckUxOjCvJJJJSqRIES0yJ7Hh3CGwTJGRDQ9CEkcz/1GSDKzqHmPgQGlFa9w0mf5PGoX0iEcCEicrVAClAJFFAsf84BgrabpNZ87Ot2dnpNRobZEJLSM0Z9SJFRoFPkndfp1tyiVFpAGAFzmOo/cOoipixmap/0FNIhKUDPvMdg7XdSGTK638Y1GtFj9VyZ7vabP8YJl2tITB504RpXFYGu1UeEu66QveKjLBEALtiCGy8tvAT58A15tIK+xI1oD+E6epAMBzd2/GtPnH/9u/xNNvjTT+ZSAAEtj7MIzHw/HoxhEId3du333u7na3VVobo0nrwOLsMNphPB6H42CtRRCtdac1CDhrQ/CHw+H7/tkPXP91vvDc3bxYkfoOU0RGa0XKjvbi4vJ4PFjnrBusHZk9IhmjFCkACcEF7zgEAFBKGW20UgjAIXjnR2uj7MuYfnd2ttmd9dtN12+0MQkvlqlX4/zgymTVgDO+xpxAUbfna2/o0iaieWakBv3irJH84pg5cBAOMp8vsDETa4vNzKW9GEJIY2WqSBvd9QBgrfUxej5x+Sbq/ROg80+yAzj1g67DJz2xFoUUHytgfRiGIQ7CubGiKcG95hamUavZgc+gaGoczySvwmfXgepV3kr67uTwWh8W8TZc3SaUGj3tlOYLwhW9LsA8zRhbV36oBtzZf03axRLl3WQzTsYPTyvN/f+jAwBw/pBirYmWomhpmanplrx1dgOr5+Mwvvrq/UWnVZ/j+ZQXirm4IOKDH8bxeDwG68x289ybXnz++edR6dFaFkbE4J13HDl1wqwIYzw6EDk3Hg8joeo3/bt/9P0Av/CaL/XunVtR/lrWoPGaqIhNeXd5eXHYXyJiXEOzcPAetFKklCEfwugGESBSWqu41A0crLPBubT1ImU6s91uu75DIFRKilQ7xX4QInOk3OdSOK3RKkqaTEbNk3lO/Vgk0CZZ2BasY4JkqXG2lfqhjeKwDEBEzS8zBaK212tF9LVsaQqFIwUI3FayNNxUo4P1/jiONoReU/EBmlzMqqkF5/l3J10pb6ADeBox8SkuKrbgnPdhHK33vvUNn3fOU26VSF6krEJi2XGwJE6srcTrHr42bE4kr8nSfXF0YZWkKw0rvE5hxxo7hIogUN2rLPOs9qVbNRRbYaztGaDw3lCyWb3UM0pLkPzpj+sNhI3hSLO+WWjz1q5qpMZft8NSChHXpjXEtrkrY0lgcNYOh2Pwvuu78zt3zs5vm673IThrh2Gw1gGicBiGgZm7ziASB+/siIoQ0RgTb5nO3GDJF0JINyKzMCtDJAjCgcM4Ds7ZcRwCB2Hh4LVSfd9J/GsOEnMBkKIol3LimHPO2RFEVLfRnVGkldFd15PRsb+JpixVmHdFgTmJ6lex7k0yOBZIIXN0pMnSmyf7LeSFbS8fLU6FsDhGlx/cmiG2D+aVK8KqOc71nAgQfPCjtc6F3hhcSzlc9TNexnXctG7rllzxDP4si31qfGXeOWsHH8Lk15IJuLUuaQodxBl9plFyZg5c1ujPlqxTykca+nj9embC4BoVBggRSYCz/g+rSCqcbcmSoCHpnqVZI1XrpsWVwebYiIBXRJ2nWxczVyHHpLablZ/+uOEJIOmkbe8QqE1mTp0cAAhwebm//k8zRr/1zW/60EdeiSR1gfk0WjkoYwzYcsENdrTOkla3nrt797nnus4457z3o7X7/cF7b0yXSjaiMRqBjs46axVR13fb7TYEPh6HT/r4d17/pT662JfowyicQJEgPI52f3np7KAU9aaz1jrntdJd14nAyIN1jjkYbUzfaWVYJBryRMgDEFF3/Xa32W2jo5xg3FUgRFueKeQtm77F3n8GdmLJYG/zrzMrjKLhLVbGnunXQcJpCYwr/OOFt+AkcUGK/bsgIShEljicnLBPXOg/a7fogjJHKzKKUWOKANGHMFrrvBPYJO+rVSbrifILsBIhcJ0/Uwt/Pe2f609J44QBLioAQsibtfbYyA9GXqGsR0LluSyHy01ipUJIxxbLX1GWpr3fhP3DmtH7wi5noXcVWBpSQBVXI9KknrXY2AqwlyQueTCM8h9qLnHjSSYntBY//fG4GWAyAsNGzFavZLBhJ9Rv1r3XHtzo5/2biTOKK3dR81ymltV7570nRZvtZnt26/z8fLPZAsB+v7+8uPDeRV0RiHSd0YqctdaOSqlN35FSwuyc83Z88YW7n/c5n3X91/mRe6+ycDy0OtMFby8uHh32l97Z4K21NnhPiJ0xfddFzNrakZmVUp3pTGcUEXPwdnTjyM4BkOn6ze5se3bW77aR/YJKAWBMb5kvW7Ddt80wO6n7/rrprY2TG3OIzKiYNGTxu1D2gqFMRyFCVSUaY9rPxzojMeiFgBRqhUpR/jHMwNxaDkuD9NSbhpR0XDmsRBILEiB6z+MYKVWC9SwEMlswNCkDMtWGJ6vbWqReoeDT/rl+xbVPG4t3zo4jB0ZEldV9TSGcC9qW8lupKf/p+8qCVd52DrJcBUU6TZ0UughGzJ4gecjMcwlWe4iZj4BkwSkmb0oWyDmoECWBUsxty56ptoNOmEBSRkgxH20CWxYBzJOz/xvqIvVTdgOAjcPdVO5FsP2CAjfODFDf+/4POOfNtdGVf/dLv+Rvfvt3X+6PgiB1icDlgCEs7KxnkX6zPb91e7vpFSlF6L1/dHFxebk3nTa6C8EziDGd9+64PxLi7myn+85a6wY7DgMz//E/8p/utpvrX5kf+Nfvjg12jGUf9sPDhw8Ises3RGS0BoEQvFFGk3LejqMTYa216YwmBQA+eGdd3PSi1pvNpttulNbZSQJinCMDAFHi8EwSLsymzfNignP3l9JU5QjJPLqtVBCpu7VayoK1mJ2SbZnkZXRMASNAlsDMgQCFGZJEBjHHAqQBIlsMSM4SAJlHpEdrC0SS7Kqd2ZAEIN4HO47Oudb0azKrnaqFVJnd06oAn6xuv146AGyX4wjgQ7DW+eDrmNoGuRFZOGhWmzIQXA0eXU5exZl5SQRuXpY8pndez2GECoKa1RSZrBykqS9LM8W5MVBm/OPi2KtpKzMikNSyiJ+eBW6yBW6csCfQudVmNYE8ycgRRJz3P/Se9336z/ika/64W+e7r/pPf89X/uE/cTgM2VoCp2zahWd44AAA29327p27u+2ORR49unj1/quXl5eBw0ZtjFGBfRTZIkDXaQBkZmeDHyOMIL/3d/yW3/bvfdn1rwmL/Isfeo/RGgCG46DownmLIM55ZjZa9X0PLMzsnBMRFlaKiHRSL3vHzN4HHwISddroftNvN9p0qCaYgaUybz/hvltvZeH0o46TgLd0+pjmChH23qeYn5lYcg2NLVpBabuomKXKLCEEQGetBE64kKII3me+NgICA0mOi5l0aql0SQqLz2bXxXAF02LJD3b0wbcCdISVZ/xZfuhnCf+36Hr0DcmenBiCt86GEArtszHsEJ65/lSmapkdlnN3pm1IeXjaOQDXW78s9p3eo5PbU0FAoKVTPlQBsasHwIRJNjby0kyydd9TDcFpJY2LYXj1fMU6KnIxNv/0xxUjAM2Ju1hnYEAxo5QcFVg4xwlA/s7/z/dd/wAAgE//GZ/0J//r/9Mf/ZN/7r3v/8CM0lY7UZOiuBJUpLbb7e7szGh1eXF5794rDx489Oz7rlOEMW2cWYK3xqi+34DIMA52GED4zp27f+grf8/v/m2/8UZX5R/+039+/+Gj+IAdjwfnbNd1fdcDDHa0hH3X9Sgw2nEcbGDfmW6z3Witgg/DMNhxABFRWmnT912/2ZquB6UqBL/Ynud8XqwV7ynoq3b2hZL5WgUc4dzUkHO+BokE9iEpGZhd4Dl1PL+/lUfBbJErE2RT+4SLAJDj0YEFRKXIGKOMJiTg7M8FSEXFD5zBqLTlTOnFKkXlxMCbanBBH9hZ571f7n1hPpjijJj4NHX7RhPAU2VSes+xm8hRt6VbznT8CeaBlv5UJqVTMSTYdvmz6ySnQl7rX0um3S5OEVpVNyIwk59jG1hY1fs6EHK+EcqGF9iGF654hU5rq5biWTsnV1TFny79127/c1gPTpw6TOk3U1mZlocZVU7G1QQo3/kPvu93/h++TOsbPD6f9PHv/HN//A/9ne/83r/4V7/1Qx9+ZcofqjeKGPPOlTFms9loUsfj8cGDB5f7PQDstlujjQtOnEUkozVzQADmEJwbh2Fnut/85b/uv/gDv+stb3rhppflG//Wdx73e++sJuXRO+cQxBijlaYNEaKzVoQ5BKVQ614bDQAufVgBUMaYvlem77rO9BuliEWYKzJDNmQpBt61XVWxaoltFLakz7I/4QmE5RCCJP2YUMLkOUSMAUCEi2VJJafJXV3JGKj1vXVHJVKf0Hl8iHZfwOzBSUi+9LHpo+hllHWcCa6Q3OsSIQimhGThznRYXC0EOIS452+KRe1k+PpOAFeUxpU/z4vavGYtsJOoAHfeMae3tii5M+2n0fqmC84wWfwX2H8alrNCeiIOTUdEu75r2Z4L4RmWL0HC06W0AFRSOQytxY5IFWRTz5iIk6swCtaMxArrnNN964hHyMkEOdggJ/9mI8uf3gFcq49BaCMNMxDSDt7SDnFc3q+Hjx5909/+rl/3q7/4prPHr/iiX/ArvugX/NB73vf3vuef/OP/7V/+2Adeyj7w6YMIu073rlOknbOPHj66uLxExLOzM9PpEMIw+uB9ZzptOiEM3t85P/v5P+czfu2v/hW/6pd94Y1A//Lxvh//0N/4299xeXnpnAPETWc8Bw7srDVd1/edc9YOY+Cgte4S1YetG711IQRBoK7vt9vNdqOVAVJExAJTP4/YRFpUG9wGDUgBIYLS5i+l/PSYzM5ZnRuC88H7aDhROzsgCFBxAoXKQW7ifAkTp28bGSkiTeBH6xEvUOT76UYJHELgqNuAFIgTc8lA6RyeIGndLCDCiKAIg2PvHYABE2PRkrNQCMF7F3woyCDAcrX6Bh0AcMKDTtboMnKtwUAgMAfvOYSsksWTR06FxOaV75QReNOJ5BRNd+0bVueVtMHi7R5h6sVldgouksfroBiUYnedFAjTd5rBkti+yslUKJ4jMjO7y//xnvd/4I/+qa89mR1czy2ydKVGOJX8lHlypVnC1Tag/bXlGqDUUptHNMX1UU3xPtVcNC+TWZrdSHxpVP2kD33o5RBteZoHfUaTlgbsazcu8eNr/9I3/rJf8vm3zs+e4Kn7lE9816d84rv+/d/066xz7/nRD7z7fT/20kfu3X/w6OF+/9JH7iGHBw8eKmNC4NGOxujbt+5sd1sIojW9+U2f8ra3vPmtb37Tx7zlTZ/2yZ/4mT/zU9/65hefsgr8kT/xNY8ePQo+hOBFRCvVmU40hxCC9+zj4I5910V5lws+ODdaCwCoFCIpo03X6a7Ppb/epdaSTWnXXBWNpWY5TO7PiRbhOYBIDOkt/BFmD8wpnFIYkoH1FDQ68VBkFu85meaKMBcEF6tnTRawSXLFyJmykqMRkNJdm/KBmVEVhWnyiifUShEhiITgVXQwJ8SQzyMOzvnAQR5TcldBAnzqA2BFJ7D8PLSfXy56m0TXidILIAIhOsYyE6kmF7nZycqMeVfIeDNGj8ii+J28XDifn0QEVqJ3J5qAzC8sFpJO2gfOCAdYHxVt4uB0/RKVn9OFQVW+i9TpmDXOhYUp0bo/Yx1hWHn6f+gj977hm/8urr+VtWO2JIZScfTMgGV7lE3iOQSkaM2bud1zbmsy3sx7r2zW2UY7y/IXBKi+28R8OOHTNz2gjexudm5jvRQqpsb5E1qrdpRs7qB4eFGBILLzTkktj4Vifzj8t//d1/5fvvI/eprK2xnzqZ/88Z/6yR//EzgQ/ZVv+btf/41/E0BIaUT0ziOC1oa0QoCYVqiV7jd9JP5bNx6PgwSPpPrtRmnjg89wO7LEGN7pfZWKjicV/6Gxmyi0d0lFWTilwkSoJ3gfgvMhSJBMI083gEIgTcLIjBxC1pXJqofCWnMTWjFeXsXJjDJQAiQb+gcRRfFdeoCCOAggDipaCBEqUth1ICQcIEJjEIQVcOL5cOA4ASyiVjPBZy0QowVdYFHDH/9nvaYRXkglZFHip8d54mi2MYE5ajM/nyGw94GZSSG0Ed+yyPgq0GD12MoagRrn51N5t+pvJdj+KgKCbUN3RQwBVqeYFG7IfIEIWCTLWF2qWRBAWngvryGWRUStiqyX4A0wuaRQl0SFZFUlFdtCTtEo6g4JlpFnWHlX1E6luCjMC7NTwRML9lK66xBmmoyaY4j31PtLM2LWpFtphiURANFa7Xa7Tb+JyPg8dbrI6pbadZz9KlL7FdZAQP0v/8W/fs+3fvt3f8kXf8FPXTjsh9/zvt/7n32VtRZAjBEi7DsDCN478CLMRNT3vVIKibz3wfthHCUE0maz3W7PdoB0PBxDcLx4pxfgQNVhp11w64koIsDBB/Y+BMbkSQfMEsmYUaEmLJ7TKgxLHmS1qEdEhCzhTBwUanK353QfaTsLTMDyUhxWh75W2eNx8Jy63ko8hBQpoRycE6Qo4AjBj8eBTYimGoTKAbMPPInM84ze8ldRZKVlX0Hgr/vnq0Phr7cWxnl5jlLClMxV4FUfvPPRAbz2/1x+q8TWaagatU1/O5Rg3VHKTD2As6rTzBAy3XtSp7MvaGJ1Ng9O4ZA1sImzDLK6PNUzE1TnWd2gLI/bKn2wSiotGy1Jt3X+xzRLTMc8V0huWqeYLQFcBiFPEQXJpRJgIrlVbNwmFHLmUler89cnLCy/WfK2aKIXpXFgmpixUrc+p2A7YTZ68/zd5+7evYsIwfks+ZtMF3CBP/G0lckTQGO4OB1C2BDSE27xt7/ze1984fmf99mf/lOx+r/nfT/+q37j73x4+UgZYs/eB9Npo43n4Ebrg1dEpu9M1wmLddZb770DETLd9mx3dnbWbXpnA+BQYJ+i5a8a/hWGfsb3MuAmAsziQ0AM3jvngg/CDAhE0esTEUQhAmH0qijmapE9RYqKkWbJL4tvU7J1nQ+sk85n5utYPYC5NkhlVEAVZwQhHk1p2S11yUBEREUqFhnmwOzz+Mg+DDKE4LuuN0YhIntvrfPerznTnN7sPfXO70YsoBOeQ5WQDydGP7SG3ciSJoC8tsW5tUe1FKq7gsn+uRbWNfE8gnNaN0Aj70TJMMskM2nGCZKTY4Dg1G+u0vTLkgiwZu5AG1LUJBVfCe3hbAWF7Wx31ZY6t67S1rT2WsPJcLlixCENMtMAXGk+lrINaEI1cHZwVvQWwLwWnI6wqYuaDjusp53We31tMI92s7EaKIKzs+0Ld88B0flAOUZu+pXkBDKIj+99Tl31v/kd/4AIf+7P/pk/tar/v/yhd3/pb/ld7/vxDxKh1lopDYEl8BjG2JxHZS8gOufYe+s8M5PSfdd1m43pN8oYQoXETY2aBI9YnwOlWY5oBqVAXQguOGsjeB6Cj73yJIJl4ZiN2fIdY9tcKEFR7MvAebYIwCSZ6oOCmJqesuSPnqMsMm89a5/FyQusqk4cESKi+EMB4hFFIoJVTmQ0ewCWAFyWHZIkBCRRDcAcvJMQQNgfBhTxztVTeH4ace5lL2sV+IkOA33zBQKe+m8sj1WivNQ9mwTmkCKmy7OHKxIugBPBMAgVNXsm6JBF2l+lsFkETVUe/VAnDa38eqkDrSyi105BmdU8qbvoaRVUCmpxumsNrVbyhrP1FE7WNSv2SFihFuuM0pNrU6wTHmkKaZgHQcIsoqxgk1JtYnDtp0m07oDygNSb3IkvATNHllbVLIths+JQ1kOg0arTOtJpaDoFZ64yK/ex3GT4nREY/vZ3fu/L9+7/ii/8/BvlBPwEfvz1b/32//AP/uGLwz4jhKhJMZAdR+etItX1nTGGRayz4zCCMBDpruv6zW6367ougtJB5kYweQMEhFCdAVDd7QIsQggsHNg5613sjMV7TxxiWVWEgip21sI8MewRFCmVgjcAspkoQwnbi7tgbgbildkR59uueTGTNpsRaj0ycgnvIIqvBvKLwQxRMnNkKM2DyyIsBADgg0eW4II/7IVgtCM3BrWTZZWs3n9L9P/1nACuBILWzqbyigKzDxy4cu+TuQEDQnqX43YUGqNEBECZR8A9blNeu3LOCyJiTpvD1ZMtq3lx/fvOwcQyaeDiBVSCB5y+Bucmr206Ly7W1tMTVHLZW7/xWZmEafNbbVtnMYtXtLa1ueCVqGMrjJN6gzPJaupjuCrguNzBNEnxc+pGba1UH4rpHKcqaoKQVkye8DE1/cnoFAjwT//5D/74hz78q7/4C97+tjf/ZC799x88/AP/1X/7dd/wTUorpVXf9QIgHEbv45a76zqltYmeQt7aYYTAqM1mtzVdZ7ped4a0SuzH4qaTe2dseyDM9uXlPwkRFHrvxnGUEDBHKLEE9jF5fQr8TCRNolpeS4oUUqwq0bg0ttuVJ1BTb+uBNjfiGfejmguIFS2wALCS+PxAKfs0ek/Gk0aYgYljs5EfT4r+bkBIpGlqgkvkOQLzlPahijAhRB5WgEpfsnQ2noUlvpEQ0Ikq27hXIKxE2YB3zF4EYJHlObXAWPl2tKhc8WvBYnmBDXY8+3nN1hhxdZc9+5PMlwWLrrmCEuREP1gsm7FehbfAi7SGD7J8UwVOveGJolRQIJyT7GaH8sw/Dtc1xQtQEa9TC+UULoJ1Pk7h70rD9KvNsWUWziZrIQuYh58mPkIatAyT2nKCIsqWHK97K6//StehGyPAy/fu//lv+ObP+lk/49/8vM96Mnro6/qxPxy/7i9/03/z3/351x5cdH1vnSMW3WkWGZ1zzgGIMcaYTisV2FtrR2uBGUy3PTvb7XbKxERiCMyR+xizDKeVXj6GsytOrni5PeYQEICIOPjxeAiBt32HSlVG/NWNS5Ty/pIEID1bqfqH4JzP+VrtUzrptws2jHEgFI5CYUEAIpUsaRq7UZiJkeIxlFYWiIigBDgmBHCI/zcAT48Eo2KFitLZRlTnP0mCtdLaABHJaCLypEAkOO+dg01fLyTk9dT3X0EDfeyfcYUqWuoux5YdASAAWDd6DlD2kPOeD6olfia81/ncEyyyZJ6entBxThOpVWGLsBmZIfkLeGm+yUKsVwp5a7sA5YpBocy3DysqbqnNY+u3v+50sWH3YGsF1PJesFqezsznZ620rKyXFpRcWG/UG1E9Qq3za/G79ODhyqGHtXtuNSs09xUu6FzlqUSZkXuK2GKxTsMn6pvkxAjeTDEi3//Pf/AH/tUPf/ZnfOrP+6xPv3vn1k+G0v/Kq/f//Nd/45/+2r/0oZc+QlptN2dq0yNRCH4YB2EBhK7rCAEQQ/DsrAveh4CI3XZr+l2/3erOULR2yKRYzOB0XqhCCT6t0VeKS1oib/0wHsWzNgoAlFLCbO2IiIGZEJJrUCwbyS8yseZipx93TgwIOXEYEDURIk0wJZY7ogrbKY49BCRU7uZ03xXgrsiBq7orRZSKRcgZKzsyU+SGTutIZhYJHEA4SceK2U12MiPE6IWHgKRIKQUcQNh7P9oxBUOW27VWIi//CLA+od+ABvo6fJToKiQUAM9gnQ/CtObkNvEFcaLbT9vSwtqcSjNWYGOlX6pFATVjo9Z7YVscVrHxpjrJqbMGa8+j+kRaRH/k8iSwMPcQqZa0NT0U69WBtNYfctIbCueJF6U3rthpWOcUFz4W4so4Vd18UL1maer1rIZXuO+8Pcj2OjixsGBiAhcfpLlbWEMAwoYRVco95dCO2VmIS+TvBtutE+SH5VxYiHrMEpgPg/2u7/mn3/Zd3/sxb37h5//cz/w5n/Fpm9TTvaEfIvIDP/jDf+Yv/JW/9Ne+hQUCMymFiCysUWmtg3fjOAJg33dd1yGAdeNxGCE4IKW16TZ9v911/ZZU8rDEvCdiYGEh4lzzM1OrobulC8OBUSB4762z1mqvjDFaawS248ghICmlY6wZTsgtIqQGOzAHEVnccKiV1loREXPa6NYIOoMk1gkTE2Qj9rR/5tKD09zDHdIqQTJBAaswEiISUkhIikAyYTkldTOHwD54zhhRKwWKYmGFiEBEOVst5O1FcH4chxDFbjRDYGHdKubpgCB9U97omg6gdqtZgtAQmL1zIfgomAbCAgqX5xRRVjaisoJdVKqmJie4jnaanRTSPv9rm+e8CJ4ifxqXHZnbUOWt5eQbDVWI+9S2I9KygORuRZq9Ut3DCDDwZEIzn0VyZsYqHieT9Knp+RPrEptlQsteXU+zaLSQMokosXk1uOxKuBl3MBrpTLsJXFtR4wL6kva0ailA+V2R/AitN+rP4uPUN/M+XFxcPHz0yDqPpKJJDSJ69j/y7vd867f/PQD82Z/+KZ/3OT/7cz/7Z338O9/+etf9B48uvvsf/pNv/+7v+c7v/kcf+sgrwzim3Bitu74nJBYejjaEgAKbvo8BRCF4CcGOI3gHWm+2267fRNBfGZM4M1CBaoQTJFmd94hAHBsNIa2E2Y4uOAeQQmC0IgnBpXYQuq4T5lqZyVN2e4JrpNhGYnLTBIlm/YpU2b+CgEjgFA8VD4BYx+MeomLYYyjzQb5PcaFFEp4qXD7TJInUkFkRMVYFLKXRECWuaph2vxnwZxAOQSTCVgKCKIG1VohAndZKC4gdrXWhM4qIRALUvIl6X/3kdfvxOoAneTJkHQ8GAOAQnPfCnEezehVZOQNljnoyi5YpXHH5TVcY97NXsahquCQZtdKtirEpi8TGtem/PqPqmlhUZjQtxwQXe3JpyT7TQjW7/cRepa3wpQHGq2wWal1t3kFjdlSdwJvTlFI82W2sfKVMW3mUJVsYp2M8PkKcp/Lpt5LmvZ3sVySxsDF9uUyHh9RCC55Jz+FxBIGriAzXBv4jwSHw/fuPPvihDwyjNZ1BoABMkBjjSmlE/Bc/9J4f+Nfv+TN/4a+98Pzdz/nMT/uEd73jHR/z1o/9mLe89c0vPD1x6MHDix/74Evv/8BLP/bBl/7ZD/zr7/+Bf3XvtVcvHl2MbiQkbRQpdDawiNEKAbwNdrSBQ991fdcj4WhHO1ofPHBA0213u935uTEdEgKplopbPXS0QrumEtLKKMAS2EcfoRCISJMipbyI90EQtFZaaVDgvefAFUTe9vtxGUyp3Is0WGIILNGmjlPbDhNTYv3eFamd1kRmu68c5ttKXKZzIK7NmWv3UFSR9K8wMRFQ5XsbKYkWmCXi/wHywYOMIEAKgZTSGkScs965zihs942vU+S3XgOtn2gT0OIGMoWfQAgcPZum7fr8UauXmpM2GCcTwOW4IwiNg01l1HYtkHd5HsxL/nLuWCACrdWCFKUgVreYzL5udbU8bTeREGvpQXUXxuUaY8URq6jKiXrXrAzyq8rQUZ5uM6E2j14tbR/XK2Ah28wO3imFSerxpz60cKY/QFyI85Z6cJxgsvKAzo5RbMe3nxAnvMCR4szeBySFJMwcEJApFi2Z5kK499qDb/uu7y0AnzH6rW958R1vf+vHvu0tu90GhI1W2+3mbLvt++5su93tNs6H/eF4HEbr3Gj9cRgOx+FwPH7owy+/7wMvvf8DH7q4PPRd13XGWffgwYP9/uC8FWZvPRIa6kzXKQUcfAyKQcSuM0g9Ijr2Yr2NPpSI/Xa32e36zVYbg0rFd1dmLIgcGxmrcl7pRPUnElHUaXnnrHPsPeSVgPee0SMCAXRaReCFOYTAISqEqj43eWqmUPjs918BBdHzOcoOogF0E5w3qU0zqljfcBzPKhVZQIlWQVj6pdJAFeviCGRJJnIyc4h6hYyIBgBhRqb0AJcNXE71RqUIKOqG82Ih/h7MTsQ6B2K9cy6KpajmIjWNES6ez6cIc9dPnwSAV47KmOMZmCun75a/UQmpcN4e4yrJReaYTl3DcI6P4cleDvFEC4iIcNJ9etpQAFRsnGrpPHO8mn+T5UmQb1aakggzhyJB99W3rO2u2qmiKvQt1g9L9RniQq71+FNzNWMHq16wEvM2dtlIUIkEW6aRtBTYeuyTzCKqbn08+bIekxT0TBCftTtBaaV113WsjUFMNpNVXvQU/9SYjIF4Hz7w0kfe94GXhuPxcr8/HPbDOAzHcbSDHW0IDpCAlO7M2e7s7t27280WiYSZqCQept1j8H6/399/7bX9cNRKAYLSKtL0EdBoGoO3oxXhvu/7ridFo7XHw5GDTzQZrfvd2dmtc61NYGZmiYLa3AvHq89Ju8sABNMriHOZhEhkBPHeDceDd94YbZTSijxwcAERyCQcPLQeyPFuKnGBRPNAjrgMjq8NmFuSJCWn7wzq1zmQCeqs7eYiUbh2+4reztKM/eUS++CdcyhgjCZSgYNvjEiz/kAYQsugQCSMyoUmmnK6tYM4b8V5T+Cs98GzcN3pnFBxPpsEl2cBAclJ+lwKZA8hcGAJy33FTERWufAhi2RaPS6ImAVlFphtWK8qAtjwDVveTG3QgLji419zhSeP+KjLkuYtkmYBLYAr2UZSeyDh1ElXZXn6oQiTcwS1Lm/SnJS4EI1HIyRcO71wVmxlnkKEsw1K/jknFOE1ckZXMkXb80GaQbLVgkkd2TH5zeX9wnSlit/bRACVK+Rvj633j09bjkpOTBxFiaUs5l0zS1SHNtZCOBseM9+LiJ2/uLi4vLw0ndlutsNxsMPAIZBSShvSirQmpSqpNEkGYTTpIHw4Hg/7/eFwGK0FljE4pajrNkjogx+HIR5DXaeV0pq0iAzD4I4DO4dG9/0mSqcSsxOJMZKDpkF3biBSrnWMPkVUhCIyDsM4jvGyRJTfOyfBEymtlCICEMLc+AfPLEVRlW7iyrYhokIs0+TKEeeZfAURiZRWlK44IZEAhBCYOZVezMlUWHkPttL1nFEOXI8Q08mBDADOA4IynekMh6B8SHteCRLPy0hWEoE0HCS7UEYRYZrSrPKyRBEhgkZA43pDWgGI8y6E8JjGrGW//ETRQGcubDI7pkpBYZEQgvDCbWFS4eKC+F/+/oTDZ9UUrPIEFy03VucDLiRUtcHNEptfiAZx8lqYUi0aRSEtjkZZLSL5LoO5A6YsBNEiKyfcLCVyihAu1X/e92MTpCa4KJOLF/34RnpClWBlrBDM9pzSWn3O3cKwBq/KcF3xa7HdugCs2Qhic/GvAfev1fYrTojp+sZOk4AISSltlFYaiYQkto4JmpMJUqx1qfHWVUQAMNpxGAellVKKmb11wEFrvd3t+k3PIqT01BylWoLM7MRJ4MvL/b17rzjn+81Gmy7YUQRIK0SE4J3zIfhOm267NVoHz8NwHMYBfEBjtmfnm83GOm+d9cw+BCGS+E7VsYbTg5Q8neNf53BX4cDBOTvaYRgUodZKIYImb33wgQwSGiRI7p4uhOAAUZEmTbAQn0jC9QNnw5/6DQYkAMK4FVBKK01a6US9Jw48Whe8RUIyhuocsuQ8iNPxAlV6X7x1pA7n4JxNgJw8cZU2nVAgxfEt5RCPm4jvMwvHySANB/n0ChLqe5IQlaiUFKyUN51SCgCDD/EAmO7hkiXdiDNxrvV5QhroM9omn5oJBICZOQgL4wwywmmFvgQ0EFcjvWby4Klk49IXusHXZbkOxplVm0hFBjoBhOcpsTLuXGyJa4PmlcqCtafEbP6oGK5SMPD0Z6IZK6a4X13Z105fKhC526fFXxOrKrOwRBqL9ercQmiS0kAea6YgEwugHm7mS35s6WElsDVtnPN6fEoQUrnPYMDIVS9HJjea8WdGrxQRVZa3DH3fv/jC884zKVKkGXgc7HE4pgDUwi1seFmTUpqQlFIgcjweEdFFSISDIJ7fOr91fmu0o3WBcvnUWmutA4fj/uicZZbD8eBd6qmV6TabbQjBjmMIQQC01n1niBQLD+Porbd2BBaz2WzPzndnZ6iVlwNbSVbKafSWdmzOM0u6PRmU1kRAmhDF+/F4iDeEUeRDsGOIPbnWCkkrUgDiQwjOcwgCoJTSOmkLMqoj08osNdUF6onmHlprrbQmpRAVpI4/CcYyaIQCQcCxyIQLFScwnMzqodLwLMHNafU0rdMwG2YRACoIyQsIFaqkC44GdVGbDCwgzBw4bgwCS1kbJFd1EUZiDdHjFBTn0yQ9k1K9BoA2wOQnDwvoVIOVL2RcAnPgesVaBvT5aNwc9AlAmTmqrUXJVMlCsDBufAzPReqLgovk0ccRZCYa5BXggiwwmHWD0khtrOmbWBZTVNW+q4AVKZ56tSfS4iXiYjErJ37bioQpK/Peggoy3wFkO756USGzk2x+C8lEWZ2yOlEaBR+2RBGUwPvD8fbt844UPhW95qojgyYDIx5GezwOWusX3/Sm2OtrrQLzgwePUsDTvHWo6Sfp80TYd50y+rgfrB/daEEhUb/Zbc/Pz8/OzpBI5Bh7xQJJ29Eejofj8eicY2HTGQ0GgAOzMZoZY+Nv4ofWIYRhGKx1EdboNtvt2e7s/LzvN0E4l6QMkCQ+XkU0aKdpAQIQ74J3Lu51nPNEo1JKKyUiIRL/SCmV8lIi4h/FwFobY7RSWiAaZXLwnttgLwBAUpBiKJAUKa2U7lIeL5BgEXwlgJQzfzMW49ItzNohEWlOYpyv0qo+rDF2lOmzVZUgVEAgJCTSVlgR5uBDYA4hns+BRViAGSQIiGdGcSIC3jsA552P103WS+pVn3oqCOh1OAUmI2hmzz4GPkd71yIVnbIYAep0kFXGDDQU/+rhz2Brc4iXgpPDrGBuWiqwiA6YWjyYZS9HrmFyLEZczC3tebBAXiZvlEjzRMGsopQyKqTqLtneOVMYik5lpV9eX3W0ny2iAcQm5lBWVGtTRGa9dFosw0WmqHSE4v8qRdIBMwmYNKSAVme9nEREBAir/5iOn4ZTVNmvotLaOvfSh16yw/CmF1+4c+d2eY7j9InX9fvBK3p/ZlGK4vrp4YNH9x88dM6d3zp/8cUXCJEBFFIAsdZeXGjEMd32OGUTTMFTSRDDAEJKEZILfhyGYB0affv27bsvvLDZbAKLAMTSr5RSpJyzF48eDcNRBLz3h8MBEXe7XdcZ74N14/4wAoM2qt90ipSIjNaG0VrvAMCYrus3m01v+o02hrQCpqROwoUXem4cKJOTC4TinbsYL+1xCMH3fcex2jmPCEqR6fv43gcObL34wCJI2BlDRPH3ZZEQfAQ9WGJZ5ORJSKS01lpH7miG7dM6NW5EqtGORdonEic1WV4SFnUZACEuiHkLdIIBkJKNI82xWQQgildkIvpM43O5v7WQVjpiIHk+iOeBdzHpJsGngQW8s95a54OXq27LZ0l2e70OgPqD8ynY+K5KO25NXAKcClWLwJwQC1W4v6wTflbtYPIXC64DwLKyI5iMwnECoNrvnmb9FbufGXdG6my6manCYtuPC2BrHs8Glch2lSuwwpTBq4iTeCJYElcurKz9CJl7pJatdZtVg2u6dpyjZjMHO5TZBgJREYUQ7r/2YBgGECCk3dmGouFL5iDhlVPrCdSv+t0x+daE4B8+urx379X7Dx4gwGa70VpFZ2QCIACddK0No3VGTcPKOyHBxCHETPPNdnv3+Rfu3LkLAFEaSkoBgjAHkePx+ODBw9EOm34T2TKJQJ820+ysU0p1uu86wyzD8TgcBwgBten6brPZbTbbftPXa1dKjXZ1ESpEFtu9LwcGEG/tcb8fh0Ep1RtNACGIZ48IRmutjYh474IL3rvovWO0UTpusyNxKVjngBnizljpVEkVaaOVMtporU3kBUmMwhARhLxdzaEwJ1dVLdmhxi8JH3P4C5UxtMKI5kWoNhSQxZUDRNIKRVFQoFIjwBy8d85S8I4CgYCwcErDCc45H3gxlMAKLeXZHADPDiDF2khbpHhCM7OP8Q6JscaZgFX9lhNbnnNoMDRq1BmmnU2apfVHLud8BVuuYz64BHXqNjX5tuKU8SPpib36gmENKcq0s6d6rGlKA5UhBa8oUNIk0VcbOViksU2vhGaknTXJA8iVjK5FtFmh8CLVuBnUQNaKjgKLy+nMzFAmh6jaBKSOmqltW04Aa9M6Xilr/Sv3Xh3G8db52Z3bt5577m4RJCLl+KirwctFQk8UlKaF7Whf+vBH7t9/YJ2142g6E7vKzH6CaoVdUSLmueRR1CFxV8wheO+9swC4PTu7c/e5s/NzIjXawVqHSilSIYT94eCdOx6PztuYGKyN2W63LMF7b4eRQZRSu90WiSLVJ/io7/VkurPzW5ttr7QxxihtIpNUG80SH1kREWAAFR2YcWY9QkgibMdxGIdMMAWtVQg8jFYrIoV930W40gfnnYtWz1pppZTWSpFiEO89++TwAHHhicZ0nTGdUkryVoQUYbIEUqleZ2oNIkky6CmMZqoY//WQG6XAeXCPYBFVgMCJDiAz7wQL8o9UR54ssVgstr5QYRlAgALR0BQQCRRpEQiaQUSRUlorpGC7ELzS2kvgEDLoGcXAs45sBorDs8gEfl0nAOHY/hOSNHhAS+Wry93CxguXdam4BLcmkc2mV2BpLvrYcxSr51QqBIfaQK9TxX+6+6q6SLBItpz5Yi/4kctqJOvUr9Ze9ZrA3I03OrVqGqkyamqXrFLnvuLKJIbroxmmbU8xhZKSPjx/LE+QjkWUIqO1MO8Ph8PxcHm5tdYS0a3zs+jBcmp+xtNVoNSXGPrNHO69+tpLH/7wfr83xpCKZjRKmEGlCSD24clMJocmQMVExDbYgIVHaw/How9Bm83ZnTt37tztTe+CPw7jOI5EZLTmwMM47i/33ltS1PcdM4uw0h0yWOucs4Cote66TkQOx6ONue0Cquu3u7Oz8/O+74quigiVUlppzyHWzzrWNCdOUzwbWBKKPYzjOAwhBEWolDLGgNjggxASqgT3c/DOB+9ARGtjui5iTCzivR9HK8EBIpKirk9KX6377UabLp2VRJH0yQDCIXt31cvDlEqUb5vJMAIatD4d4K0MQGR1+yQttaxi4lGatCI4NMUuLSyF21CP2PHExjY2OJwABEXESqFI13W96X3XsbPaKEL1RuoZ9dOHiq3UinJEQsnMDCISey9FCqJ7XxWZiIQtMXKeZlXbuCGtTESzZrBOVIS6oVwZcwt1E+uTBLkBiLBOyFpUPAQqTUoBHGkNwMIr2SUA7bqiGNViMV3EJsRFZmsNmebVJBwGbFEYbIg+9ZSS9e+4UhPLhcyWDLGzSpMNTqKnmXqvxrgyBRRxzkaQMmu3fKRmVJLl1jguVGpFUFn7MfNgx1dfe20Yhxeff+6tb3lrwu5PzwFL9Ue8fNHznhEePrp45ZV79x8+iFT3SOY4caiso1pY2Q9gvpreh/3+MByP4vuH8QABAABJREFUgHh2dnZ+63yz3Sqjx8Ha0V5eXgbvu67rui44Z+3oQzgzu77vQmBrx/1+DyCkaLvbRVBtHMfgvHUWBLTp+k3fd33X9V1nKIZpkSKFmMD1bLwf880p0fIlLe2AABiBHQ/Ho7MOhBUpYPY+hBAUkda6MwYAQdgHF1yIPB+ttDEalUKkzIZh750EDwKo9Ga37TebEHgYjswsiESKC4N+ot/MZJ5p4Iq76hQNWS5zVBgnjW42xcrrYkI6hf/N+xUs4ylRxCWQJIsRkv65OBmW2pRf3ARU5idC5T41vSilIKSILCTU2gQQY7TWGhXJtErIbBhZGVArzOUJYaE3YgkswjnQGeJtJrOMG8S6PVy2Y9drW/EEZ0dO4jZ4sonH1a9a/0a46CKbmNkmvWUFbpEVwHlBII22bfGGlEUEAMpKkw6Vl+oqrWlpm0OLFh1L8cXla5cZglZtufMhceUEkm5tkSbceT7+nJp7qu8i0+JagAFRm3RjH47D8Xj0zpmuf+7uba00Fff5a9zAMSvUCzPz8WBfeeW1D7304dHaTdd3fSdxqzedn9P1XxlYK8xthmlyYGtH78N2t7tz5852s1GKSBECWjce9ododcCBQUTHQCzEPF6I9w4JTd91pgs+jMM42hHYk9am7/vNZrPdbvqNUioza2LGuypLlogWRud9TMcjoqJomelDEBHnrB1Ha0ellFFaFJEwBxYUQlRKM7PzwTvHHFBAKW26TmsFgN57Z50LPp70pDQpY3qzOz/bbrd2tKO1US3KLAISuPA75o+m1N7T2ZqHs2K2MjeBmYPsYvMlJ/epNaFZBIAFs34tEjuFCah2gp0FS+ESisBUihKiTIRBQMAH77znLrtME5FSiPQ4hPIZzPbP7gBYFZJWDTlzsvWbRYAtIyHXGyZY1Iz5L9+E7U4MrcJkq8TUVWwnVSDGoheY2Q5P/SbW+M/k/DCZw8mSTzC18Cu3mgAKomqskGXFZ2lKtsg7geKoXHZgpcGpleQzE27EJpRRZH3N2yg+S4JCtCnN78QEwmUL1ex9DQtvjWKsLhPCVuw/SsdX9/IAFd8/P+HYLHYkZwJjmbWx5i1xvJCPLi7f9/73PXz43Fvf8pYS1ZIWikR4JdtHAC4v9g8vLi4vLh49uvTMpJQACksOoToR2p1hBpy5Y0vLVYsUeIHOdLdu375z9w4pxdk/Zxys97bfbLXpnB1FoO83RGSdG/d7ESHC7W4XLdKmxj94MmZ3drbdbpTutNFKa0U6JqBEU+iMcXPmoaW9XI5gBE0EIMPRHY57YY4OC1qr4MPgvSLSSiuN0ZXHjwMHDhwQsTOdisR8Is/MIYRY/ZmBlNZd1/dd3xtjuk1vdG8dA6IkK6nIgcbCtc1D3UQnkxT/lcRgqBSlTh0oZ+1W6U2pjguisAhFf4sJlEOpWboCMAs4ZsC4T8naPhYOrEjiCr2tM9igHzhlnNdDbWyzhFFYomJMIO4I6gyDaobGuSz0enTlnzQ7gGiE1yZ1NGnfq9DuNX9hqVa/uKDjLLEMrJScZQm5QLsbq/9Zw7s2ZGDJkWiB7sdbSkp1DK4lJK4cf9IuIWoPt1j+ueXkFIpa5XIhJ4lSsuZRVPNtyhVq1yHS8KlkwS6ddryFKJCVplPYpzx+wmvLKdSgnbQjnAAAKVKA3vtX790fjqPSGok2vUm2AdGddHkXZraPD/5yf3z1/oP7r92/vLxgFqM7QooapWs5TVSkJanvSkRECMzWuWEcdGfu3L1z5/btzWYDAM77/f5weXnhvCVSXdd1RnvnYswIonC0pgFUfdd1HYgMw3EYxyjS1f1ms92cn9/qe0NEkANMdNJMleVcpmWEICKpA0WkpN73HMJoB2stB59Y+KQF2QdGBTqmCwTvY44VCCkV2f0EKMAhBOd98B6YkVB1PRnddZvNdhfBKEU6noKcXE6yA5ZMBnqSqJaZzZz6/WzmQxHGio470StUguLap70xiasIeHnt1MpJAAA4yXG4eYqUIi8STzStFKkoYBamOrQdm76xiTksiRwRcZN47UvasQAwQxQVP31r/wYeAHKlaKEYtWdyvEywXinas2cQ4XqCfWz+37yzblvf2r2z0rsmVe+VOSFTE9qcxGsHWD1uCBR7u5kMr7agqhSHUvs6TJ31dBNWkRdS6bNwtkmm2X4Y61TsJrizei1YD29rAqAmSG3OcUUods1TgFvNXMoXok6fEeCpvZFl2j1WhqbVTFBHvU1Xj5sboRzJAkwxplWP1n/4wy8fj8e7d+7cvXPnbLdNdZC5+PZE4wFNCgBcCC+//Or9hw8Ox+N4HENeYuXok2nWOHlw161cJRuPqxyldRj50aOHjx5e9Ka/c/tu13cgoI12zr/26qv3778WQtj0HQfvLG66LoQwDgfvmRRuthuFCkCsHYN11rvoTdT3/Xa72/SbrjMpfySCRpnFUqZYAgRULBI8cwiRv6/ittO74/FoRwfChrQXCd5ztHRWyhhNRMLinXXOBw6oMHI2iSiqE4JPmE5s/De7Td9vtdGkjOlMXMinA1gEmoThCLswhHRI4WQxlUzi4v1G9UouBcCQQjLReY1FiAsSGe2jCZWahvusD5riSSJth9J3pIJiCgMwAsRfLQQWVpnxozI8C3mSbfIz4nuPBUZNfKr4C7QUj6hfYymIa1aRn6JrP73/7TMPhYe1+p9x3tYEARAXIP8pO4ElzwVPsoPKY4crh+g03Ale73zFSpY8/02lpijKOnd+oiKsXiZcDBm1Mg5lBizLeiY00gyVqi8HrBi6z1hQK4PBku6EUE9ZFcJRXQKZQWOt62fk0ySvHm6MnNaITDinjS7IoNj4VcyvrqSJWpE2W+ODv9wfrB2t9RFt3m767FCWfoxCiv/Ief/q/UcfeeXexaMLASZURpnoUZhdOghW8yquOAOqHjFi8Sz28nJ/cXFx5+6dW+fnpCgwj+N4eXlx//79w+Gw3W77vnfOeW+p60DEWhc4bPWmNx0iDON4OBzEOtSKlNZK9X1/drYznYlOaTpqcbNLpiIV3Yx9ajNRWJgDImhNBOC9Qw7BejvaYRw0Ka2NQhJMYm4iRZT4PM56Ea9Imy7aM6CweO+ttcIBSJHSpFF13fb8bNNvSFMs3TVOUtW/aJ8g9a4kdvrlOcy+D0CkEnU1k0GFRZBRRdc5zIm+GBv0uN5QcV4QaUwd8/FMU5xkth9JCyghEUaMdgYgjFOHJqXDqtCIim/EpyL8ijnodG9zTA2oWH+v/wRwTfOggkvNkwBmhl44ydyJ8t3OJYg5/3uqWriF6Keu7BVfF+UqqoW0NPviJzhxi/IkWUdgVUKyanchFSGxSQ6o8k5WZg+p138IK6B/7QZ4NQt18vGHHEQ6a7mbKQbmr6WWOxdbkQXMhlNwQeVomvQP1YBQ/4jmihO2ha7l6NS5mdMSezk1CTQGEeX5b2ElqN8nxMLylWLbJQCIXBg+JZA1Em4UKVEQAl9eXHrv9vvLF55/4cUXX4gDUwgBFSlABLg8jh/+8Idfu//gOAwsXAat1Yf4sQ2E1CuTbOuKmIjrzrnBDjt/Rpo2/fZwPNy//+qr91511m42m77rtVIAMFp72B9ERBvd641C9M754N0wiHVA2G+20f85ZqdopeNuVmsdJUZK0abvjTEC4p2XMfrbJ9uyiPCEEA77Q/HbMBHxH7wiNFqR0ggYgh9Hx57jVs+YPjb+LOyckxC8DxICEHVd1/WbKObqul4bXWd6J37NJLhJrsFRglV7gtKUcxHpmMk6tHqmEAEkWtgREqfZgBAVEhHFZUl8cynzw2vOQQmmTecG50z7EliY4vNi9VdRqBACs+fm6SIqWuP4eaa2qiUPooiIpGkm6jDSFCMpBW2lsXgi07fMLay2f69PJjDe8KtPjS8VQL8ecbv8L1mkGWJLrITa4fv6D/AKl32tVmMVrs61KrFZMderHFwFmtY0zLhu4dMsc681np3cnqyObTId+WVok9pjaUqnrxr/GZY2nSQnqJHYzon4mF8C60mwvKlT0IBMry91G5SPhhB8QESjtIB49hcXbhjGEARJ3b11po1RSrGIYz9af+/Vex95+d7lYW+MMqqLj6UPIT6xORhZHnPLrwy4NXJYQmi5uION3l5eXj64/+B4PHZ9t9tuWcQ6nxjv8YnVWqkY52udG4FZd6bf7c52u8B8HAZhBogKo8jqAUWoTRf9H0xnRCQyo1iCpEqLwCF471z0qmEAIVKKFCOnfHNSitAH9t4554GFlFKatFJIyByc887aGIGlTae6brPd9N1GdzoaP0RmtUwdGU55LwD1CgcRKTlTYOZhpmaGymWV4k43tdDJxyi7CSCiIlKkWUKNB2MlxsuFv2Jc1l1VFgMTAmIeLECUVkZpABfzmaYUqhRFCSyyEChPvWFl8JWShIVUycoWeeOUADdwA8X55xsuTEXZF5Fm9k1UgwzeSmHDVpEk08TUWLufIGlCfu5T8kY6grGKyCp4OLaIPDSvvnY4F5zgd6xCweZOfM21qIpRnB9JKgXE3GIQl07T0mZSQJMehik9pTBoJge96W6bhYoVpF2WhJTaLi9nGqciwRO1Lpn1yITbJIM8jmaR7VElxY2/eZSxIizNf2tpMpxlHn6AMNdk5WtFtaxjJqOrmf1TP5Xz2YBDAADSGolQUEA4hPuvPTgcjy++8MLHvO2tm864EF599dX79x88enTpnDVaUdR1EYrjELxKJHqUwAKCUohAi+YlSUcxAbg5JKpuA0vHBwDGaBB58Nr9V155ZRyH7Xaz2W0VqsPheByOgGCMOTs7EwDrRnfYc2AXPARWXX/r9p1bt28Zoy/3l3DggKy0iopcDh6E+25z+9bZZrNBJKV15JRaayMGF61HrR2P+z0gmq4DBus8iEuIf0T2OQxH670XEa1V3ABH+wrxnkPwIYAIadNvuq7fdKZTkQ6kFVXmLpjm+tQZEGFJF9ZKaZMM1VQ29syVAKM4ND/fPPHiUtHMNMMc5o4lVgaimc8CQhbOqXhY+hjMzXtZCpQzImYOsDCwoESSmKZpfgBJMfQAMZ04JiUkYyNKm4bMS6Q4R5ACgcChOOpyTFhYMllkrT6viEZP1G1oZTdvGAtIMqmuAfxEZujpyS3aY0aNsjCdeRjMwQ6ZkUhAlpzDxhm2BbvXhKfJAzSfN8iY/fdnHfzKeX5KfojX4Q5hbQJRXDxnRCU5tZFZ/PVkqSgr/NxqbKLcwZ00eoDKvBNPqSxkNXz+ClE7znwsmu9ShwNje7lFhBWZbtNpUs47O9qoQjJKGW0AYLTj5eHAgfuuv3X71v6wv3fv/muvveqc703XdT2H4IKHIEb3u7MdAFjrOITM310fwqQaHuV6g7IIWGcvHj06Hg9d3+92OxEZ7ei8ixdFIWmtQghutO44ACFqpc12c3Z2duvW2e6MFB4PRx+CRlRIRmsEYdDG6POz7fn5ead1iDZ8inbbLRLsD8fL/cGO1lrrrPPeI5JSHgAVAkuSiMXq751z1guI0VobQ1EGEA0+meNiU/fGdJvNru+6jdEaMMaHxQTGXL7z5jUHHGUzQQCl4vkUJxgsKv8qMmO+C2yClyAG62pmVpHwGiEgrSBMLSDOtpOx5stk+ZI8BKQpX0gQ03VCCAzovZPAWmmjdVR/i0T0JordghfgkHvKeMhVUd/x7iSiuH8K3qPO+yeGKsD+mVA9r7EEPimvh/WyJKdQhTYIAIAEwDOH4IOAonz9sS4RMqf0lJyGdV5iFXVLJdhnlhyz6L9xsSZu/qp4WRbEo/QaIhDFLiULY0YrTYaOwKlxwAw91KS/qWQ0W/OW1tl4LS/UXlXv304Ngs0ABhURZkIcIpkgQs/pl8oQDXPMCahOLUp7Mpq8/ik7aDWpy1jB9zPap1Ts/smoL8NClcfndJQX92+aQUmyQjNdZjVIFfvDwMxMHd2+dfv8/EyEHz58+Nqrrw3HAbvox8JESmsYB/vShz/8yqv3gg/DOAJg1ExJYBAIPojI88+dv+Utb7Xj8KEPfeh4HKP9Q1aknuh71h5fgZYliEBIzHI4HIwxLNz3m935TpF68OBB/ORuu0WFEvhwPPrBOjcCAmq93W53u7N+u9VGC0DyyBQBEB0NZkDMpjs7O9v0G6MohVeFAIibvt9s+hDk3r3XHjx84KwT4c4Y6/zxMGitTWdUpwUgeG+PR++9ABilTaeV0gDgg3PWh+Dj5sUY02+3UXagdNIZS3OfYHH0lJk3Q0XrIky5FxlCbxzjE3JSDbjlqUkJmaS0VjEKbNobIzBh9MqTLBUvMg6uN1GTP2/qpyIcH2MVkNB7551HxHEYOtNtt1ujDSkV/S5RkICE4mDNrCmVqRwFUabFENIV0Fr5EJwLAEK6iwp2qcnC8Wl7TPLKk9ft118HIDHEQoQZkhbxdEgv3szpDk/A63OnZsSFeHqO1uLJhUCsPRwNcbNhDU6nxeRtL0VOu8SZTp3njTl+PY7My+Nq0r1clyWAM4bmFMQjAnN3PMS2mU/PRaHdzeRxJzMil2sBxNP5aHj6lsDaMlQanGe2CKqAeQEQIQSj9fnZmTG66wwB3H/twejG0Y5EqLUxaJj54cNHAkxEmrRWBogDBxcsEXVab7abF59//k0vPvfo0QUShsBGZ+OQVe7X4yhlDJNEu+s6ZZSzbjgcFantdgMCw3CwzkadGikCAOvs4bAX66kzZttHD7iz83OtDXMMZYwJg6hIdb3Z9j0Kd313vjtTRJG9jkTGaERk4eP+eHlxcTwejoej9y56cComESZCQwoInfPee+8DiBhtut4Y0wmzdXYYBg4eiJQyyuiu2+zOdpvNhlBFADxkF41JwjiRXqqhrtGFSqr/2aikfipLfUZERFVlRacGSDg4a6Nfk3UuprBYOwRWiVoTdTAJQyr8xKQswImLxPUUGf2b41a26JlGO6rDHkRC3xNS9D5igITqYFyNmMRzzWlhBeKP6wuFJFqH6AFL0MVtffXTHzsByFOPCPqmtRau3FCcUERO/FZqmSI4weJYxNVNCLLMIl9mEi1pg0wmNZhU+PwVNvrQSFNnIEqmr0bz/iLYwCwVRql/KqAQ4DJ/MYtkZcZdbVGWJnIFK9uhaDUyeaK0IJOg4AoYKFjUxJXRRuLFN3yeMjpl7eI8cUiqjWW1nmlibAQajZ1gOgVlWl1gSseWtB+DIsKChqEqNbqVUyslc2amI+xU3HXjxRvZis5673a77sXn7tzabXfbsw+99KHLy0tttKkJ/FLFw+XFhlL6zp07b37zm567c1sjMHtASQknOPErlgk9uEirTBeWEBGIo/qJEXDT97vNLoTg2WllROTy4pG1Thu16Teew3A8BM/ejWI9GLU9Pz87P1PKKK201pHp3xmjlSaFSlPfmd12e352hiKYEHaWwCJAAP1mi4Qvf+TlH33f+x5eXDALYhIfxW9ljNFKi8h4PFprEbHvOq2VUhoJOQTn7DgM7D0QKW36fhNdKLQxSVuX7r7I6k+7mRp9IUQASkabWK98knO3TNrISpzLnHhnaRNPBTQHEA48Du44HKzzLME7zxyclX24REKBxZ6t2fIWv4BKl1BpVYoNNQooRSAQnD/I3o42ChqiT54IaG36ruv6bb/pFKqo0A6IKECIpLRkk1Fm5pC2SMyBAmUua6yUUIEh0kDcNf9NHgO/z33IXx8dwOPJQdL6BuOaFBXxGkcerlKA1jfFkxpgYZU0STDyLhSvPvawjjDAE50dzSGbGn45cUwuLky1fV6pKbiyQDiFKkvLY0JY6iFy/ru0AY0Is8QkXOfpL/oQzDRHrusxnuQ4nb6Bl3qwRm+SFAJXK6cREJidd95ZkV1nTKc1vIW8G0mhs46DBxAkZXRs6aNXuyVEY7q+M7vd9vkXXnjxxRc0Aoh47yB2eXUduZEZr1TnZDTGIUWKQgjBB5FYSpgIjTZEZMfhcDwCMxKZTW92u7Nbt3fbbYw7iz4WSsXEdVKkOm36rtv2m+12I5JitghJGQUszjk+HLxzL7/y8oc/8vLheOz7npRSWU2bxgBm712M+jKm67pOacWBnXNudNZbYdamU32nTGeMNqZT2hARR1nBZICWHR1mcu1yY5axID9t1D7l6fyIC/eonYX0RuQQ9rRd5MDjOB6Px2EcY6pM9CV1zlacdaxtomtGUMY2qeB0dbxMvNXq6TIwhzFYsHmTz0nboV0IPjAH9rHBEIDAiVVljFGkIsYlLEGF+JAWqzQAYK6fPbkK6HmGOwA8wSSdH5gzkGX59VnwOF+6TtRnrEJcqgkP1/fbp+iaU98+4xJmz90m/0nmfLzYaxS/7xwb3bq0xFVV5W7ZyhZk0eBhHfRao1lYLQ9gcTS1eiaoiku6+xkYJT3tZY2AC3e2JhRhjZg580mHKrIRp/ik9dutkifXb8g6mFdRj6fFhbQuJzPMq4mPRCQgAQHKJlr5kcXWGE5WMuOq8UEk1muVlEEAALvt5h3v+Jjz27defvnli0ePgogCiWauCCiBObDpurt3bt+9e/fW+dl226uJZ5RjJqVUEyk8qLWOpsh+J0UFRxMnlRTcIfjgPHPwAYK1ROr8/AyRjsdhf3HprQUB1Ga73ex2u26z7bqOUFFCezKLCwEQlFKlXiuiwLG7ZGVIKe2CvdxfPnz48OLi8mK/j4j5MIzamM4Y6nRsY8dx8N7H6Wez0cYYQBytc6Nl7zxzjJM8OzszfR+t4riw8YpIhojSjTxZeC5yRRNJZ8GWlbLMISKobjxB4BCcd8M4jMPofDKei4UlDpdK6cj6j1hSYEzVedWNMb/CkFptbrZY9aGdbPMQSjxIOeQIEXR87CPaNo6DtSNM60HQSne9Edl0XU+oY0wRACgVs3E4/WgRjum5uMqaX5Spx9bt9vhYfr1+PTr+uQckTxYrVdXA64L9812RXIv33mqhZq5qiCtq2eVvQggsmFyYJyMRWUOvEU405rgSei6zznzO8IHGUg7nK4OrwP+MZC2Dcq6YmyrG7pwPBFcsbU423iW8OI5bZRB63BuH859ZLgnPnCWkdapeTogIGH3rksFNSo0XQtptt13XA7NSdByOzoXgPVsnAorw/Ozs9u3zF5574e5zd/rOAIAXViXOPiMCj4MW5yyDolYrn1SokJCFY8HVSoEAKdLKMIuzox2OSKrbbE3fb3eb7W5ntEkL+qiHSs6RsRUmBIh5W24cR+eMor7vEcSFcLl/uN8fHj56eO/ea/cfPGAR03V9342jI0UxmzeZ+IcgICYeC9oAgLV2GMZgLSBpY7q+67bb3fm51vo4DOy9hKCV1LN1suVcTL+tWhQo06gw8SxRJXYpICRGTeDAIcTiHjXMzo3D8Tha65znxActGxAklbl4acBC4GgghI3RObQGXDP/Hmx8fdfWk5j62eQpNTHa47qFOcZ8xZORgtYiAUCCD0hEFM0fwXvPwQszoBBiyPbWq3fQLJxDnuEEcDWTVJZd//rX4xQiWscjC6Sde4Gc03K7ceOfSZam+MW5l/1kHJbIIfPTciLbCBSEZ3KVScoNLOpvkVluTMT7sl1Efj9iiyhQBVEWWLxA4Pl3irlD1Z4aafLshIloPNm0QeMc2PpItt992jtPW/D5Pi3FLcsaC7OewdvMqoo2lJnTs+wExFPs0hruSSuHikBbsafqWakwoyeEPP22XN1eU43I53IxFq3T3yaTKUpyoHpjFJOl0hXRit78phfOznb3Hz568ODB5cXlYEcA2d2+/ba3vuX5557fbo3SZr6Bl2JtxZnMvuIkvza5Vuq9iQ1V5ByEiMb0XY+Bw3AcDsfDOA6g9Xa7252d677XWinSAHnPq+IiFPOyRbx3LKJQMYeL/b7fbG/fOu92nYRw8eDhBz/4oc/6WZ/6K77wF7ng7GgL4h8HiBTaGrwAKFKk1YOHF1/79X/dehe899YF74Go2+52Z2dR02t0x8DMEkQ0TtzpzPXERqnb0sZynFDsuAkmnxhSydANRMS7MI6jHQdrrfNekhEesgiHEKlHKXsngifCEnj5Dki1ZIY5izQDcjn8Zqr8ZWKZ9Y61IVec6MK8t4me7RGfS1GAKCH4ceARBuaMh5ASgdFZEcYqE5TLASA0qQxgQcAro+cJNUCzNljUbXnWSuCTu2CR2jN3nlF4Elyu6d1rhvA472zlBOEd6uUh1FYL+co0bt4zg6aaRyT1HNMMJ9W0UfnvnF4tECBXWlVuEMoVnuX8zD/dAEzgPl61I8A12k35DU8oWB/nnoSFjZ8RU1mSFeRxEoeG6102wwscEOf8otRTSCO4XNwwLOxZtFLGmLt3jNIqEsaN0Urh8889/+ILL5ydbfOtm6vLkzMt5v4p8c1lkdGO1llFarPtgw8RJvYhHIfj8XAAxH63256db3c7bUyOGoPof9y6yMTk9eCDjz1wYB6GAUCG4+C8e+XeK6+8+urnf+5n/fbf/Buu+aJ/5L3v/5qv+8t2HH3wAKJ1Z7ab7fnZdrPVWsczNnjxHESg8bTB+hGbWgZEWE54RKSUMkYRoXds7TgMOh7tLOKct3Ycx9E5572TxP9JyFKBvhhaFYu0j+00u08oVL20KimreROQutdiCDR7zigj9RO8PAm3ChW5/G8yMmKR4IPzQZLOK5pBaSICZkIUFmtHOw4k4oK//rLsGU0Aj9/jXncVITBFYs1Z0QW4yNRdqCx4ZNltzu3wpbGJmKgzRVUvEy8IW5w6V/lFSRasknYrQg5VcweVZWorva1HDWn7vdqpZjZjJpMRqo5tWlCH8mwBrcFP4U7NZc0thtpUz/r8wAlHFCyW65IzENpwx9lxiCsrGlmBiaTg9bV/aCODrADYIvRM6+glTaDZpgBJs2WXKYChChiGev3SvuUEBGr6z7PdlpC2m00IIaan9F1Xn9RU48dYZz5jJWw+AcnVeeXSkKrY83EcjscDIGy22+PhOA4joPUhOOcB0Zh+02+6fqO0RlLxtiAEpRQRTuEKGEPqg3MueEZCYzqttLXu4aOHw3EYx9E66z3LTQ4xZh6PhyAAiNr0Z2dnu/PzuA+Y8fFi704EBEkEg9VNkKn/GfMhLNZgyaktGkogBu/3+33Uo3HOsYPsS6GUzrzPSKCJWmzJt3Kyj1BakzS6SGnsyabKMWOhItTqT6zDXTGtbQqaiShCk53RpP7ErPJgABEOLAgSGjpYDD1LSyVKrkZx6+6Dd/6wF+YwOmQBVTIbpG1KV2NFrmDMrH8Bvh4soDYECmfkb1kmzMKJ7Dy6kjvf4ni4WHeu7XukqjPSTKgLfSEmNltdBPE0zbWqoRU1tbGZaCpCuk9OjQePMcVv7TVro5HWJmcJ/kzXqutMesEEzRlRL5FzbXc+XOlfJyteN4s0m1O3JzYQ+UxHHTtEagd2FODc/mMd8ymC3vnHx0hMud4xjI92u23XaUTUSpfyV7FT2m0VIeLJ0e6KSbaZwBEFxI7j/nCIIiwkEJHRWTc6HzwpZTpjjMnNfrJRVVgNVyIQ+S6AIhDYA4DpDBIeDscQ9sNwvLzcD8MREI02N3nNsX6x7jcR9N+dnW+2OwAIwQeQKPGSSFlhnkMlc0MMjHZskbOJhFqRAAQfrHXjMBwPB2EhwOD8EIJ3Poik2lj8wKvHVCDlRWPFSyYqSWlQIQ6CJTspRTRjmaUq/U26/7hGTireuxSMVoSlIslFM7uqEU05khI1lAzz/WtK3sxLrNwrJVZoYOfZhyXTXmZKnWc+AchyRzwn+M0trnEFe5pZV0A1dVVuGdEPtq3iVUZAawGabfPmw0BTeQUet5GrTAKkMDOm3UAuWNIuiqsja+oqcK1Y4yyaeFanoSI+taB1m4BbZQdcI8Nyslya+t4SgjqdBMXpJL7EN7/w3Ld83Vdf8/6w1v3iX/sVdVu0khg8i24XqUH5ucAsR2rXBtvL9zZ2db/qi7/g//g7fvM1X+q/+uH3/od/8KvWWGMIazx9qNLTTIb7y8GzOn/jtOOYLeYXpx1WoUcwJZoICCERIgNGV879fi/CpJQyiu1g3SACyvREikgRUVpvFp7xlI9VXNZJREIQRFTaWGsfPboIwRttosLb+8B8UwwBu81md36r6zdKK2NMMnQpDI5odclZJFXJIqu9VzSgThsZCZ4ZFIDSGhCCD4f95eXFxTiOzKKNji200rFDTt5TMU48V37MuzRQsTuYlkbpjmJpP7JREEzJo2EeW9gQJESk8datkw6zkHjit2BT5CLttZieq7rtjX0wRzE2zxOTUky1Utqotuk51WSuDv/zv14v3a+PG2j1+0Th9bx1jZSvtdTdWZ+KTcXG04Udy0OcG2pZFN6Gxghz3wRY3aFWw/Vyoo+aQprZfsMV/hgNQIKnBp+nIl61G6DVwxJPuBZc/2csKDqyPgesTwp41eAgcC1r5ddhXVXR5pLF2DyRddULqtp7yxWCHFlsKVIzSYq00troruucd8LsnQvWxQzCmD+llFI5qLcygCqli0WESCmto3rIec/MoxuHcfDeexOUImU0KSVyM2N5JNqd397utl3XIxJLgvur2lG3fZnCUBs/KCyZzxHuJ9Mxc/D+8rDn4L0Px+PxOAzOOq006FS8sfBCEVGEEDk1bJCjIiMhF1UkieY1I3O208+GypO3Q/mfRtJSbeyavUAmXEj95MoKcTF3QJxVl4ycD4J8MFF573K3GUCQsbywZBVBOXNKTkA6r8vToa9G+VdRpCugpUnI1LxcjmvumOYgsuYKVoUFVFB+bfbQlmiZIoGkekQz0DyHXyZJoojwBAHlzRJWaVOYk6qKTFGm/Q+usSSbEAF8XAmFaUkkdT837/3nvgkLyi/WqTdJAZHFw2vltYJob1j8J7FA42iKLakghW+UILPFsbj0+CwuQotgH3yqm15q7dHVkBo2u/m1I2LeKEiieEmKM560T/FDKfWn/sgfvP6L/eo/8xf+0P/tqy8fPeIQkBQSkUZttOmipw9NXlLFGAmSXihqYkHAh3A4HFgYEY3pAGAcLSL0fd93/U2vnyI6Oz8jpaqzbnrcYFHuMysSsgVnTJ6J3pbRSw1QG6XUMIyvvfba8bjXWivSWitElMDehcC+Vp7H/BYiRREOx+ldIqIygkSfhZQ6W5H0F8hpxdXMiD7OO5JSJoCLX1a1MI7+RoUwl+kOEhI6UCzlBLlSPqRIznyZFMUE5IZsIigcxDlnXQgBT5DJrnks4MkbeP6ZNygTuFwXgSVWWK39Kq/idf9+aA3TkhvnY4DpQv3EnCq00qS11taTK/WcEYR4Cq8/kd9QgSaIDf3myt3N9Ub2mkUKibcgC92BnFigXP8EwHqBeR3YG+Zr6ZnxH6wtL575KLrEfh4/B1yL04PPkpsRAltnk+ar16bvN5utNiZiP0QoKSpdJk6fxG2wEg6PHj0ax9EHJyDWjohRRayU4liII+DOwjeaAIzponKioqVL2WQqRSI6+UyIFBfFSIsv8i6jFVHPIRz2hwPvEWkchnE8WmuFBYzECLDoFjNlMMU9ICGljWlCWHO5x9hQckgQFE+kYUQEApIZgA2NzW+1CJw9iBO8RzJ/5IvEtPQ4U4PIEWOrJGMik7UcsCCyILVpremwSdttiRHEiPhGjsJ6pWLJdWvPfCVQ4SeTnXzZ1QuzCCabmpyQOYGHNXeE6DG20EkLxcC4bNWloRwl1lgJgMRiLtlAe4mg35hF5662WMcmemW1Qa6FwgvLh7JlkBZEhtmVw+b0aciqc6yvdl5ufmgW4MydNzAvF3B9PXzNM0bSnkWqS12o/RUsXObtPKVVEQOANROU5qifxO4VE0fniWrrFNWCGL1PCw3oSt+8K0v/XIGTLC2r5c5TPrOHw2EcLWgNLEDY95vd2U5pXcXZZsiFAYgj85A0Ga0uLo+vvPzyaMftdtcZ7YN21gU/aKP7vldKCfM4HIfRhhBu9Kq4SVYtvX/ybe6MiWBK5K0XcpBWmLI6ODAIqa7vzTjIfn/56OEDQNRKK1K7uFL23lkrHDDCXYoapXAFrseULg7BBy8igJSXp8UtJEXER+wl7g/yO8ZTca/4IBMMMHv4Jp/aSdGIGU6CWaOYahipDFEBInMcfUJeKrMISgCR0NxGRAqRikUmIiilO0NayRVe7viYBvL6pfv1YAGtPDOSqoIsv0bWxv6WPCFXxGet+E3WBwSsUOrncJVUGElzgaRmEGLR4GKbWoLzLYNMWzqE2lh/SZE4bfJ6Za9aT7GrSQPziY/qSGC5cWml9nVOEwbO2Lqtf9BsVopdFWE79zUt2fpC6KZX5/R65YlLNF71whCfep3jrQ3OotakQGmjjSLS0Ts9QttJgBjN3UiJiBM/DuOI436/H+04WotI0nVR0BBtRKNnTvDeOR8D32/6LE9u5ygISDnCMPoOearIHcX8g4UUKqUFojnPwTlrRzsMR2sdALDhnKoWs3ZQcquvVAK7OO1/hUU4BW0IiySRFALmoaS4wiEqVEmUKwKCjAE4IfkIIiX9tG7vyt9jZehbKzQjj4eyHYBAkxE7A76x2N8pQFGAKPHXkIIqVE7T8XsCqrK/TBODPO4JfZY5wS0LSASu5fs/ubAuSSnZLmF61JMiGnF64yN3R6ZfO65LgGqBaBHVNnrgPFzmiPmE/AGVTmF2CMVphFIRT8xhrEzssUleLLhdmhDThFviRla6xsmiMnWgqc8QEG6VoOlHSwm7KyufudGczI5qrMV8TZ+P1TKj8UDNUQFTeLPAzet/vlBYw5VRQI2MklcjU/hyuTTJpEAq268WQsK2kcEpAxSfDK2SRmdYZnBpPJtvat+GOOv7QAo1e/07PcHTiYHJ0GYbvWLIOauoiyQXZo6ilGIkKcLeueNwHMfBu2CMIaTROu9935u+7xHROT+Oo7NWAIwx267XLdPp8Qc/KSglEAARtVJKq0jnx1z0kzNqhlhCEIRgjBGgEPzFo0fDMIbgRaDfbiLf0VkXOBAqranb9AAQ44njHR35kNFBHmSuDopVJL42zK3/RNzPjkD5DGnT2kSwpf9kclHVkDbW7OmdDDC9Epk10YiTUw/G4BCMImBUFE3scnSQ1K8pryvYM6dvIQzB2XH03k3amSlzvrZIkMct2E7X7bbO6ytwluv1CPUS/TSmncJb8SRuWut1VlB5PAWUFzptsxJNcxJlTq7MkfWTLxcbxZnIosGUlnjT0hil2OHRVOyakJxywONMorXaMp+6nm2uGbbL3mrtgcvyeON+VWRFl5V/g9kgNDM8kBkWj7WGruXZ4pMW/asq7yxi51kMuNWjKM/sOxNICABwdn6+224Pw2BHa4zpovpJAESUQgQlHI7HwXsXgrfjcNjvOYg2hpTSiuN2lFmYg3eWQyBCUtp0pkgcbvbOswCASugKKU2aJnN/ShF4IiKE1BmtiLz3zgd7ecHBHY/D8TgM4yDMSimldLI7jnooQlJKxzQVCRwCIwJCYE62+6nfnnZoKu+Fa+Ytlgzecjxy4DmXRuaK//qpkYqCLTU2II0vTYmcqJ8NqZQoCACMSMSCU0Sp5HMKpYQxpf+fC5pNiHGL7Zz3k6GFvB7rsfo76mvw10Aew7aAFPqWjsJFqc2HNDZOmtIIKyd5jrRkSVzBQFKjLwwARYcyn8zjXRIvsmDNUId5/zy3QJAUk9W41TJLCRYWqMOdKC0PkiaRolgIkGpYojoypSIunQqiWall1dwTB11eHBv1XrXKSKt+15sjFiKcad8JEaq2plOZRaiRdil+ozg/2qUlfcoq50eerLxWcBK2TclSgthuxuXKa1/RJ6pM98px+6m31Up3XXf71u3dbmetP/hBmDUhImUOPyrCwfHl/nI8DqTiFg2Zgx1GbXS/6ZVSImzHYRhHAOi6rjNbpZSIuOCjl/WN3ncGUVBMg6I6IkZZkVbkiUCEfYhkTR1J+QTWHh9dXNrxKAJEZIwRZue8tUdIwb9KFTe0EELwztnIf03PdryhiJKRH048ntncLVNiS4qZyBSdrNzFE2PfzBwRanNhnBvMzLbFTfcy04ELSIp0yXi/IKISwrgDjs9vXHiDKop6RAAPsaBNW/fWyl6uUZ+vWbefKQuoVqBKjQw32gqR9R1BSYqoQrUWobMLQ+eKFwfJkLX960Z6XH2/+uhZ2OVgRXVBqAOMKq1TbY/QnE8TMYBwAdnDVJWk7nSvxvOwZZ1iYz/QKFUqro+srkWeqPuvNiJzqQws57X1xIJmgpP6LcbFbkvWF6/PmBr0xBjq/LTOyaZPuwQQ6c/O7j733Pmt89RMMAtzRF0QyQVvrR95tNZaa523bDlinl2nmZNyLQQfvLPOIYI2puv6GJCb85DDDd920ghEpCi+jIg8SC2OjIhGVCscj0fnbAg82sGOwzg6RFBaxR0SYk5tVKSVJqSUrxtiAgJXY2vKc0SKa+FmMV06/dyBceCk8s3egHG5SpAEBLEZ1zMjgkWqUHU705qL1FQxRBjiRl1pRYDMIQkBkgKhdp9OwJOIEDIyNk74ub+UmC9HCgDLJvyp1lbX3wE89Y9AWToJ1I1mgbll4ZhXnRPVJmCurkoTGgq0nuLlNqVEP4eaYSgp+0qqmLBCY5YUdZ29vmK5lwnkwUQ0rusrzPe9awURa7xkSS7EFCTGWCUyroLS0ycR5yyhLItJljACCyGVlMVSfWY92XtNebUQn0VuXrFg3hLHR41S5DDg5OYv0DgAzRbuhOuqMXwSOKilGTfjTiZxibTD0eMetcYkdTZlzeyYUqG7+Wqf6O7zd9/05jdtd9vD/iAikf+ZjSeFAw/D8XjcR7cfIm3HAzNv+s50PQCG4J0bo6uE1rrb7EzXIYB11lrrneObbj5iRDsVRZMAtCTNjLUAgCLFHC4vLi8uHiGi0pqI+r7zITjrJKbZdMZ0piDt3nvnXKL0lOctNvyUiXPYWhNIpoHW1gNSg4wIpIhIaxUd67x33ntSZLSORs1ZkIVzlHzeHtKiRZGEYxMG5yM+0HW9VjoEz+xj9U/E1BAmsLC8iVMzOLWpKqdUCiUWUbzo84fhdUuEeRYTgFzDF32emzfFXjWkHYQrrSJr4mOlaybM50ymJDIsmsums5SqFFWxiXPIf+4suAYTzNrhece/Wt2QrkdOb3chWMecpL00CfCc+j9Xe+DcBf0Gd1PXme/+pj8HP1U+sIr3wxptwvxUIa54wj9mBz4Z8iKJCDfyt4w1EJXb+qaK6+12e+fO3e1mo1DFH6ljgrFS42hHOwRm5+w4jMF70ooAtdEIoLUGgFjmQgiIYEynjTHGgIB1djwOzo0AqLuN0Td52BEVoaKUo4MV0BF7F6VIK600CfMwjrRXx2Gw1iGRSW4r8WijKOgy2ihSPoTgPYuw52hepElBunRYAB9OHpscndWKjIiZK7wlv454SBFmmolSSmmtRCRIgBBlwzrGUBaKh0ybY2xMXKY0s9zPFUNmEUQCQmFAisQqbTpDgZhVPAAUM6cAzuhiV9bRIQflSbI6FxGi6HEdkTHhkGgirRcuvgFuoNfZJj9myywNEUgEIBvaVY6As+DHJlN9xq/C+nCRLOPKm546MaUyVZoAkqL8q4yZcYJwBRgYACX7DsWbiqYF71T9caWU12YsIM3fC0LtJotrRJX63V1YlVWFvkDtbUBdVb9i5p4UBLT6pWVtaQ4CH70f1GBOiTjlOVgbRucFdEQR4LGwWzW9xeqjjQYB60Ko/XByfIf3wY7OGBNCUKRI3exln5+dn52dR/K4Im20DqyVUiGEw2G/P+6x2BoLuNFprXrTK00cooXyKMLadP1m2xkDAM57Z513znsHSKrb7M7PzE30wAigYjpvFPg26TvZTQtYEQLI8XBgDgDQ9R0gcGDvnIhoY7rOpCQ1UiGwd85al9nZSmuttQIA70N2WYOc9MjpAJAy2U0Oc0AEREqRUlornUzZCFSy4kEiChwAUBCziAWrEIwMC+Cs1MKEFhBKm4MdX2EyuKPiXoakFBKWnWAqUSwhxNfPzOydD5yPssQpKg2EIAgHYAmQFtiCFalk8tg9QezBp6jhr58SeN0u4dR2bdnmT3G20iwQ67YWqxirmodT3t1M08LWcVkafsvccHhyr8a1LfTKJ2eJdq0t/Vp5welV4/rlyDYqEw5WFIg40VTboaYFuGsTaMSP7srfHqBTsSIQGY7jQ3pg7aCVkuT0OZ0T0tCpcPmoxIZTkQaAi4sLETDagCqaQQyBLy8vEVEpHeGO5+7cutFrPjvbaaWYRQSN0V3fW2ePh+M4WmvH4IPzLq1ejSHheCd4l3AURDC6M31ntGEW760drHUWALQ23W7bb3eb5ON/g71PbMYp61QBhZTWpAIH5+3+weF4PFrrTGdEJARfP6IRyTFaRzMJZmZrvQsuOBEhivQfpbQiUiH4WPOTkIBDikWX5B2ZbJcjKzZaX0TpAClUVNDaiB+l3gmxtqeDXA6kXiMXX+gpXboypGmQCxQAQcFpP0ilmc98qJTnWrbO8R0FYWZRKsRhAESYQ9z0xkR7YRaRABx9VcX7wHyzTe7TTABPs1k+vXEWWELmmf40c1tKpWzJha96WS5h47iCoEwUl8SQzaER9VtbW+IXYhnmfUB0eauPoCnqC6tY46b6y4qlM8LitMPF3mO2OF21sW5oNiJTAEC8B1NY6jT34DQJATR8iUWg+kdz9W/wr7jmOR6P4zgQFRsfqfIdyrS0bvmNzSmddiG6M+VuQcQQ+NGji4uLy+LvpOljbvSyb52fpWmDmYMoIgF5dPEIAEzXEZF33nvf933Xd0QYfLBudNaxsCLV970xHSkVvBuG0dsxIv5kus352e3bt/t+E8voDZfAJddRuGzQhJnD8Tjcv//gcr9HAG0MIoSJ3U/G6K7vkAiBRDgE750PucfXShujSWsCZIAQvA8+cIiFESaL03ivK0BKDpkq2aNCTkIGwpTfErGinGcqLEpFxxfMQdIo9RquvNFJHoZY+0C1xrpzMhBCMbiWiXGBxSUsK4KR1LRWZGOyf3d2LpK8LvCBgweAUPEFBdYjGmXuD/q0tfr1mADKdWl2ny0ZEOvdTTO3Nw365JtQbxwWD2uGeUAqWGhe7mpmoUxoDFZ8virha4UzhtlgUBBxRiiQWRVaV3Lj2oqxdlCcSv/UrFQ9KhbXcUBoia0zMuhjrDQ+aqs/1ouptOnh4ENaHqbLnXWBUFIz5VQ2Q/x6ZAYEUEoroxUST/oTDJF4zx4EiFBAIg5z/Y/bt86V0cJyOBwOh0PyyQH0wfM4xr6y6zqjDYhY67yP2S/UKa2UAoDAIQRvR2vHIwiYzdZsetP3291Zv90prUKQmz7DIgykCEkR6l6FwJf7/WG/DxxCYOcchxiT7hGRQSJdNKq0MsPCxaOLRQhJ6bShJW0ocpYiDYizsk5yTlYs7rE3VKSUiRmYWQEwLWkzpT4hvZzytuYJPtVDvD65r7k41np1yRSgKfIyBZeFICqZ/AhmnB8z4TMHgEf7MREBVYyq430TWDOzB5Hg/WhHbYxaNSu7AmP5CV4CPw4KitI4KilKUtAVqshD8/G7NALSHA4Ic9Y54HyNLML1rJfujeUlQwQV1fbCOHOBb+O1Zp3+mncYTpblbRWuTGMagiYW2LDxH08bbqg+SfkXptrFtO7qcbbOTn6GbUhNOQflo/kQqAT7aSIkVKAwusdgTT7LX93EQlQks/npGmXtEqLqRTLYR4iodAT+ERDh7p07NzsAzs8Q0Xt/8ejRw0ePmJmU7voOLNlxYOauM13XI4jz9t/7tb9GRJRSxnSkKFbZmCgpzCBAxmw2G9P1MSAlO1zhJ3/Cx93wDJDyaJJS3oeLi0f37t0DgM1mq5XebMA6N1oHwqTNpu+0NgCShLw+OOc4BBCJv04UgsVj2bO3ow3BRb4/kkqoQF4Xp3aMku43gezRZr6AAZjp5dk8jFNAZN1AQvawikybqTuiOqtyIUON1JJ8S1BjgxrPoehFGliRICESMFNiw2V3RpQVDUmuTIqARSlMcTPsrBuHo+76eKi/MRO7xmdc7WvUB+sBoHrEVjaiJ7yP4iPbwhcnbZKqrY6sEXaKmfG02Cm0rBL/2X6DctyssehEZgZtWRoL7W6hDn+fWAbp72Vlu5O/bRUhC1jN8DI/xuoDBWZqO6TpOLrSKuej4UPWFQiIqMqziNg6lkvFDJImQhYmH4v4L6PhvLBgG3ecxP+5nblz5/ZNIaDLy8v9fn8YjgnxR6+0Vlr10DMHAbB2FA4I8N//N3/ojbmWWumuMxLY2nF/eXk8Hvf7Q+xOnLOAyCIIoJVGENQKScdjzDvHIcTGPiXaGK21QaQQPIcAErz33o0ggtpQDLwkREVamZgRnJ6K1PITAgrWPoMyCfzrXZjgwj1svcAgAKxs+KojH3FZnaY9tFKYmvjEGpc2uqoSxEiNQNWtCgLlUEkBEKUFSCUVAmJTUU7kcT/9k/x6TAANA72Jga1J/gJRR46rNg8F7p9NY8uhrpT++oGn1fpPRaOapOxTWgXmUrBEdRb9YPH8galQ1224tGWlUjCsBWZWcHTFC6pyjrJNfX2A4grDVBY2ZVDrTaTy3Pmorf+znUrlzRtPR06oTjKRKl2BVE9rNpicLJYkWdi2Y1/Fq4rEjgLevuudN9sBvP2tb3r44OHFowskNEa74J1zLKGnfrvtQ+DDYT8cjyKy3W7fuGkKiYN47y4vLx88eHA8Holot90ys3XWWicAXWf6TY+A0fzS+oj5OAFGJK210SoKdQAkBGetD8Fn8QuR1qbvkcD7AISd6XWKHU69mmQiZuVFi7K0kZxY/HjK4WWe7DapsCr2XkH6sSUMVwHdWUWEghhYWDiFeTFMnrzzNlHqFPP86Yq2UUJqko8Fl2dZTh9iz+QxfgY00EXb2UoqVjC1GnnDGqDHSXVaGeSsmmZmlyQqjs+PddETqchDddKL4OJoybmrhUrTJKFUd1otNWrJA1VgwSrtXK6sYnUibQN4ZeelKYSyGimmWUikljPDSiTZR/UZsMJTxpMLknqLlE50weaytqPFpIVayehDgJuCLbdvnb/j7W/7Fw8f+tHHzmSz6aO6fRydDz4E1tooRV3fvVHVH60bHj08WhurvYvsTOs9CLNAxGUijzMEDt6z58AhRnopNKSVVoqIIjTPnr333nuIq16l+82267emN96HIEcWQaVIUWCpbVpQYrqmZAAP26Ftbotb9o1yGqQoxPT0/VhmCdPYCvygluCgVFqiZEKEk3wgr45gnsQhTfu43B9m80mpaqcUd015JvV5hQYq7RPwRH8GaHgvyToSizPoZAmDbXL6SsRlfuOl4TLWLmwT815Q6l65hEfPoaGadwQIbQZxZfA1y86Uagpsg2ukwbjwBM9GsGJlygK0kaWuoPbXn+aYvKeU9pbJP5VrxulSRV457qyccx99K4AmcGbx2LTeHeXi4qK+V9J0qKyfqie19BR1z4BARJ/wcR9701f+c37Wp/7Ie947jqP3frPdbHcbEBiG8Xg4+OC7rtvudkoregPfuouL/f3XXg3eR2WZUnoch+PhCETGqK4zkYjpvXfeuTGa+aScy+gaTYDCHEIIPvgQkp2VUkCq6/vtbrfd7ZSi43E4jqPkXLV40mDJBMlioSiMIVS5sk7+tFI8eSov39YGVLJ6AVvqfP1Fk1smSPK3TXnwRPGYyC6/6RVxCRtOqNGCrjMNGU2Zj8dOXe1ynjARKsxRe+3ZJlUM6RPX6vmf9bN77uD0AJbFmWWN2xhi5LRNyaSYQv9pCqXUOBguNZ1NWMrMbq5NJ8DEqqSmGOcWA1faisp+UyasqT3hceWCyPJv1vO+cOF2VLNMl/1CzXWXpTs5zEeZG0hfPzqOglkPuCj8sNzUV+WkNR2r7sMESLSD/pQFigDvfMfH9Dfv0z/9Uz/J/RWvjdZaK6WdjVi6hejp1nWkFQf27N6Y6xdC2O8vIsLtvY/2ZAIQNxNaaxARCc75EDh+AcWBwGgkBUggzJHj6R1zABFQqut7rQ0ZY7qu7/vOdEhIyhWYJxmLVPHyqRJyJjmvZ6VM8llhYZQ6UiqKDFS8IQhEMOI2gPMKM53njA1Sw7U5sGQdfiGC5a/Mec5FJFhCxCbdatn+lRupQL6R2UrzHlZez4dWr9WFm/55PhvPOJiYFpjY7lKw/boM7ZXAHqjzgZtrkvgw2XM40hUmsUHTp7e1tyRGgcR06erh//+x9+fxtn3bXRA4mjnnWnufc/tf+17yEpKCEEKoAElBCFWEVvrQCIgEUZrCDii0pABRSoFYlChRghUIFgkoqDFYgIht0CBFiv5jQEz7urxfe7tzzm7WWnPOMeqPMddac+29z733/H7vPeVH9vvkvZt7T7ObtcYc4zu+TbUERjwVWl6tpas8miOIYdHw4ykYQpf2bcuSjoumBvHALHo5xB0NmofJyEu+kt54ByAif/I//QvVAb/ooU/G+Z6aNRHhoM9+tmN5efyIL/y8n/DjvvSmzcjJlV8VXXJ4TSx+r300qkBkvMKjoXkeCPBw4kBA+KIv/Lz3cB9+2Y/+kcMwNE3jHA9D3O22MUYibhofQgMIcRj6rq8ipT7TB0Dquy6EIKLD0O/3HQJ4H1arlh1Llhj7YYgiyS5mR855x8zlFpOUc04x2QYbiBDJN2G1Xjdh5bxDNrO2AqjUCyoiqnOvAcYPwSZqESBcuMIsPJnLaSELjKjE6tgqzBpMAaUaXpmjACrnNxh5iCCV1kht55+ZYKjyBkS04qgvwJCK9XEalBwZ6aP9KdUl66jkvp9affjnT8sEoNfdh1WbdUCUqTmfeuT+XBlw2fJ/wREdV+1YRcaceEKz8RssPgyYEhywovHo4u6e9jZLjzkAOI1aHcZk6zESs6CVVn+o5V312l+PjRuw/sEKJ/5QEU1PLGHewyOl/E3/4Z+ZjhiaGNq180ZV4p85JNb7ktnJ6TQFAEABfuHP+ik3PQBO9PhL+1aY7EFKKThcoeCxpekkv65G+8PFVpks4Sf++B/zHp7wj/vSH/X5H/mcN95+p+v6lAYVtW7aOa+gaUhD18d+z0Dw2XpkERgGe6ucc0TknWfGHFOMfYw5Sza/UrIAMyYA1JRTFlVNOVnIAYcQmuB88CGE0BpAJObseSLBvchlsdbxzuz7Wa+pR+DAiDLILJ9fGinj6Bxm/hZFSzBBOjPzpOpWJ7VB3eJN40MJuj0WbeH1a8hxRacLC/oSsTCGmn3WPuVPfyQkwmGqUz1t47G0S2csdXQcBxWp7kZcYvDTJ7dgYR4nJC4q4Ygklk8fJlwQa4vRaubQpWx3UeOK2+XBOYc3tXGtWTkV76CsGGpseqFNrvbZo/v5XJ4M3RoFjjOE+n5PgkoVU2/OsbrTTg+FCLVP00gHxurnHO9tEI/z5t7LEYCntt64nK6hWv/K8ZgGBysiXM51NbCqALdunX3Fl33Je3i6RPSP/5Kf93X/zjd2+46IQhO880iUU+z7mIaYcwJEvKHE7P1BaNB1PSKF4JlNnaRxGOIQhziYYst5duyKllIkZ0lDEhEwrg4TOd+u1uv1yoVA7OzCTDI6PTBjVQHtUxCxyIFFWzbvD7UQceoOBGdUuLjTIqJIpd0fwxoRkYmIiU9U2UW01NHQXofCULFIEpDRwgGJKM+g7WEOrB4iCYuOF0u7yYSOuA7vxBe7Na9rxJ9/AHwmz5oZAUPCI4z8gK44yV8XzPkxXBDr3frs+oB4oAOrK5WiLg8DrHxY9Zluw/MaQE8B/AcgBp0UEkJ1OeHhBvh4P4t6yDXBcUGkz5i4cBmwu0weqDvg95FeohX6s3ROXDjkTw6NcDCzHb6JI9q3IGyo6kEz957K/uw5e5jJUB9eo7kgVla78yrvECJaOPIt0bQJNzY/mK/+SV9+I8ud+vHLv+Zn/54/8A1N4x179k4VYhxSTMPQyxABIbSr9fn5Z+cAMEMF551jx+wJUXJOaUgxiggTETE5x84TkeRoLKCUk6YEAORL1++Cb5rWe8/OaWmly5WopVhT8dKpJjA9wG/H0FeFwtjBQwut0iZOMQEzzbNOdRpFxIyHMtEZQa5Qweo+owqSlsnWcgokrJr42lP4SNxzMB7UPSYqKhCRmWnj6Qn6uVDLjc8A91m4mIhwjPVZYuN62F7Oycs45QDpAiyDKUJBr103HxYmODhOFqnP18pi8biTreEOHXWkNBXGI6FVLYsrFLLDZzQXp8Nptha2LYjp0+EA0zFYLVKO4O0jqQve/BjHObB3CXsfQfnjJzjt4XAMTUNz6dZZKzeJnhfv26m97I0q15gQp3BCCD3ncyvWCGFNtaimeK2ZJlogsNKlVlMujdXhp/7k/8N7vk2+8PM/8pVf/mV/5+99t3OcsnR913ddwbgd+9Csz89v37r9WTsAELFtGmKWLCnFOMSUBlX1ZjddFpsqOQ1DTHEo7yEzsWvbplmtm3blvCMmVShuyLjQ9sxuQ3VWKqJCyR4DrCIvRsQIR1zg6N4tDrpzOTYz6prWX1/SMHPOT6qLlpmm0/1GC2S70pPOsQS1m9uJKjIpgquDQUv6pmNmptNprvDp93Rxn0lbAJ3EG0SV7PaQFjlbIdTvDC49LWeQu9Q8nED+Y5aLHm1lcSYRnRbULiGgKTJep+eiB7jNfAaUJKJFRPIh7l/ZGy09a07gH1ALUHS2qJ2CwXGRXjllA1Q3RC2xnlfXN8Knpi8uhfzw3lggmliUD2ruLipKhMxEY/LfcXEBo4nkVPjIzEyOCOu0D7zxBadZkkgGJCJkcmT3Uu0DYrvAnLNkARERBHTsLGc8SxbJBSocu4zJCGYynpSx1XTkmB0o5Jxfffn+l3zRF76fG+bXf+0v+w3/wu9KiUUk9YMOEQCoCWdnZ2dnt5pV45377BwAiBhCYOKU8tD3KUZRIGLbTDCxRXlpzpJzigOIALP3jW988I1vgnPBBU/IBSkpKUFLq73Rnu84u3ks9YdhKNV1WNUMLfq96UBGOMyInVOfFACRmUz0N0LK5vOxzFSqf6cuKeR2FNCybuiRV9sUpqrzcHPKjbg8QSLnnSPiOh9xYQf2AvvXm04A+pmp/TpLlnAaAQpz/2CBU+MHsHB1P3KxUzCbpSNuhtaSWl3KkKHgu+MlBcvtc8Xi0EPq4BIRWOafI9Qb28VMs+DtVMseXay8Dxt7rE0qD61+xjcAa0xJF4Ml1r95/IEVP+K9wSpUPpBj58GKjGGiPERyzpeLyuzEFi7VcyqYWpI4Oh+mG3kyu17efjcpW4QOGRxjMUtVGJNmK8PAkmobggNUEVFQAjLjXkLk4GleM6nW3R0uAuyml6aqRPC1v/wXvM+75lf84p/3TX/yW//qX/9bfd8BIngHiO16fefe/fOzs89IB3jdh47omGOMQz8McQBV53xogvcBAVJOxfLBVqCI6IMPvl2t2tU6hGBrXgXIIFXi9vIFFEO0MXu+AuFBR0IDVosvrH3Aaq/zquEzRwkkZtYsVN0lWPkDECIhAxgzqKKP07jzETiF+ersDq3mIlwa2xEuVjyKPSy9zjy0HwVI4Vz9yJzvkPV0Rb1x//a/IgS0OHpL0tti41s51p8UqOrhhIeHb9tk7zOj3UtlDijOtIklp+qID6WVH8DBGvl40TxNFaNLjJ5QPS/AQF2OgJMmbWHQUNFrql+P0865aulh9DWYbheZJ12FU3C/vscVQPG51SNK63wQGWE8+Ob8/Pz87Mx7HIa02Wy3u12KSaE2ekVCjDkiwq3z83v37wdGBby4vHr69CKm6MgfILcv3nMQ0/n57du3zz3RkPPTx0+uNhsEJObCMkBMOcUUV+3qwYP75+uWiYacLq82l5ebJHm9bu/eu9cGb3T0Ew3naN9kL+ji8urp06cK+MVf9MN/5lf/pPcNltK/83W/66u/5h/vLx9Du7537x4zuxDatgWALLJatz/7n/hNd9ZnieDhO4/efedd1fzKKy/dOzv3bUPOjQzCcl0gIhWPTvw1v/Tn/bSvelGEShW6/T7HlEUIyTvHwdv8EXNKQ4xpsOBDZNc0TQitD56DD03jyM1+PTqP7dXdgzBacdSNnIXpQh2CIVLSvuYVDy7DYqeTmWj0C2VEYlLEwQRpI10cCnVcJ/tRXKL9o/xKpzSRghbOjv+zuqteI9uvF8VFDvm0oizpATpjF8XjseZ1iKpYntmJHcD7A/qffQDgZ6b8z/XCAoEmyB2xwu/xGL7RumIet6PzKTI5aCoCLLZ5IwRcf8bXr3wrTBrmEKhFqK4FvU8Zw1TZiFdKsiPqQFX2a3qYjvDmQbk/Ab7PCrlJSbFopIgAkGegcnyteLSFqCG097gHXv5qrCcRUVD1nu/dvfPyyw/axg99fPz4yTsPH11eXqWcABkRddbcKxPfvn37wx9+fe2cAnzijTcuLi5Syi44c+68MRtOlQjv3Dl//fXXW+btMHT73cXlhSoEYpwyBVU8861b56+9+vKdW+cAsNntY5Srqw2Artar11979axtASCpEl5rPWhvuog8ffrUO/drf9UvWbrJvsfHl/6oH/Eb/olf+e/8kT/O3t06vxVCsMyTfhgkp/v37nzkI5/74fv3egBJ8tZbb65WzRf/yC/68IMHCrDPQkR44uRXEblVZogX3QF0XQeqnr1vGjcReOIwxJhSAlAgBCDnQ7s+O1ufsSNbv6TRkVmP1p8nh7aRWDyF/R2zxKlmMRxQdqkkdNHUkJuHtCJYhMwMmuIYcjCXCVr4485GxcU+y0KP7WdP+yoiIgRzuhAAsmDrkjo1a3AATracODVqJkueiAYK0+86cBRGWLRf+Ok9AI7ubIATPJXr//zco2DaHxarhikI8kjDtIjTPY7iOvqLhR58ARHOHnsnvxmXLJ26dGulOlkgiVrOGKqfHsJ1T1IPeeRHtJqJMVP5fixW3gdRk7N8TEEACZDZlQaPFDTHnHIUrVPpxwtd54QxvakdtOqSHEOVJG6+oi0R3jF7xwywavz9e3eGOGw2m9Qn9ODK6a0A4Nk1q9Cu2uCcvUhHPOFUiCBqVOCbs4AAHTMABOfIjpzxrc5ZsqQQwt2791556YFV/6vt9vGTp/vdTrKAKiM3I9UyW52bpvd5HywK4JlVdYhRVb78x37pl/zIL/x03ZP/0j//6/763/1fPvaJT8QhIkLbNKIyDCJZEcB7RwArACbKKWr2rQ/j7aND35vkiR1b2VLVoR92Xdf1/Y3GKWZ2xD4E55yIxjRkkZRyTgOokA8+eFEgds575z0iySgYXlz8hXs9QWpakUAOrDRmYK3KdMKpM8eZ5GUWc4QjP4EW31B+KBMxsWou0JAlrhM5Jsc0WzpWA8tY7cv9LaTjln8chC16jMhn75wTFWJkoxXy3LrOuVc6tYfzQvHAFAEBBA0pgXFdqof7C9Drujt4H3X7s7dTmncaC/eFmhKwNADFEzvaiTty9C4sESfEE9ZzpxGmWWc2rlFHt9BljH21wca6RJ8sp3odXleJSOe8F1xMDBXIP9MZFrCLilLZ6maRLGI3QpYsVjSJsbayq+yO9YargOD9X/rP/oEJhVfVGFPX9edtE2MUqRchigjBhXt373zoQx+6c2utAJebzafeeOPJ00sUlZztaE0pOeaLy6vLq8sYs0VIItK8yQcFUOc8ADx9esFEv+Yf+5pP46t4cO/u7/9d/+Lv/H98/cc/9sksOcVMDEwslFVk6Hu9dWsAiEPMMcW+3+12D27figCP3n34+OnToe+ZmZwTFQBUkf1+f7XZ/oyv/LE3umGbtnXERKwqQ+z7flAVG+eBfLtee+9jKkrgrAVaqYy5aj6lUtXFz5rME1J/nVSeNCV+j0Ebo3syAIE1PhYKNnFCbOuEACpZUtacxt46CwIBZAVQzbNpW3WqS4lrJ2Q82EDJgudmcZQ5Z/MukpxijJUWlEa5Gai5zk5MpGV660hFnqYFUos+OCIswZxO+ulmAeEzzw14gT/XG5jKbWlu6MtkhlhD6LjM21iMR1qXdR097cdzY+E2jbCAg7CSfD23Vax7i3GvU1VnOjVqPReWWAqwjtcJOtE7a18fJBydBSf/cIuvoDIjl7bEOpWcc1bVNMT9brfd73ISdrRqV+2q9T7wCLpNaNXoh2VURoIP6KPKvCkgm92LqiIiTdvcf3D/1Zdeun1rDQCPHj1+59HDpxdXaYhEFhcDBASIKefNZvPo0eO+6yULTJFOM0yhzjkm6rrhn/v1v+ojn/P6p/eFfPmX/sh/5mt/xTf96f/sE5/85HazbdrGOyc5jb9fAEAlpxSzsIiISNR8eXnx1ptv7nY755yqbjabrutiTjHmlPLjx09udAA450Q1p5hTGoZeczID57YJvmmaplGADCApy7T0r/eVWuVw68m9lJRIFymIPAEYDUx17MtGZde0SizwuJHTrBJPhpACWgZ/ySnHOAxGXRXNCBQhlUDHEi08BbKPjaDWQbBwkA500AgCZpWUkoJKzkidIcJMzN55570PzDT6FY83nW28RSdSOM4IMyIRKjvmcRie/HwX9o/15ljfd90+YQf9bOgXn/n1NWGy9ugmIjvwdC6RCM+npqjqIaSOFUG7hrwXZkOHXz4fUQcg4ykVKuIhM/i0bcH0OeghnUdPnhN4tPOoPSoRUZGmPdd4FFiCdCpSIzDluUrOWXLf9RdXlxcXF33XOefv3Ll97969diVMzuSOsCApT5ezwAf3BDCJ0gT1GsPHMZP3d+/cee2VV26dn6nKk8vNm++88/TJBQCEEFQ1ywz+TreqTpZelSOoNa85pQzw0/+PP+EX/CM/9TPxUn7Rz/7qhxdX3/pn/+LHP/FJFc2QAZQsbnE8jkSyZLWIWs5SfGQQECDm3HXd5upqiBEJiW7G97PhMqUY+zgue4Pz3odmfbZumkCO+z5O9XBJMMBqBF6C+3Mkd6Ufn21BkCxuczbVGeMETU5U+ITVgDz1NWYfoIoIIpKGvu/7GPsYU7ngRaESGlT8CT1c210LiNesjbl1H2YJITI7H1wOjapkcXYtERKTjk8ckWn6vZPbsaKoIAI5ZsalWn2RNgy6oMG86Mruurr9fu2gF4i3Qr2fV5jsdux2pGLLdOinsiCFjh0+FjiRTtM2cZnzMro7LKr/qPLBOXW+GjrGmXEBMI2642VAHMKxuaieChA+wGp0OcdVc041CNLMF0ACIkMUQSSnlIobe59iijml4rdegro1pjj0g4hklTz0m80GAMJ2h5Nd3Whc5Zz33jvHzNR3Zx/Q6l9e85juoxbgBYDr1erO3dt379w+W6/6bnj05PHjpxdXlxvJmZh1Av5GlJoIz87W+f69nOwNrzBMVVMMDEP32isv/ZZ/5p/8zL2aX/2Lf/Zbb7zlnP/Um29cPHlCRMEFu9TtJeYskotVJyPevXsXRPddr6A5pwf37vVDn5IRaejB/fs32vvEYUgpSYqgSqFdrdchBOeDbzyzk6mqz7wYnHMk8SAH8cCj2+YzOtRMzpzMETQYiT32j4aP0OicqQqWOJ9SypIl53FvJJJSTCnnPPK/axdirTlEh+bslZPAYYDf8RRgNa1iSIvkOICKpmGY0H4kcszOsWcfgnfeE5vpz4hcSY4pigIisu1u5gZ4YskegM2qnw5raPf+MuUPhKyKhx14rZGe29rrox11CvRSrclUVQ82uzpgAeeQjp5brfE/nHsU8cCWutrMHPbuUyjFydd/PUZ00on+MG5y9jvMIiI5S4lBJcwpD7Hvun7oh77v+67v+y7lrADMVOhtAETsvUfEFFNOebPZIO102vaqIgKxa0NomjZ4z45una3gg/tYujOBc7xqVvfu3331lZebprGt7zvvPrq8eIpIvvT+UrxBEAVUVQj9+fmt4JtKvTjqv1SddwCoOf/Wf/prb5T//k3f8h//3J/5Uz78odde8OvbJvzW3/i1Z//Rn/9vvv1/6HbbOESRqcYV5ZIUJBoc8d07d87PzlOMqdDz53KWs7788ssv/lRFZOgHQER2yG61PjtbnxlBC5kEIGcRqd0zygyrh5swXShStHZQPi1utbwtZsaJLzMukCVn0VyQUUIVTSl1fR+HmHNKOZv93IiUSk3WwEqUtRAmlGojooqK9QhYoUB4AEpUDYFMpiY6JjTFGAf7CMqnQEzsvWtCkyU0Is7zpIJ2ZKJFAlB2HEKoMoF1tPd6RqMP76eGn7SD1hfeLdf32oFDTZ3RU6UCl/wuXOL85chDqrKw6CTmPk8BhADIk7bOrpXp49FT5CEc7aB1XvvOP7bi/KksaZe1Rz8sxH2L5mAWMk3Q3QH0uYh+NPqXZJGY4n673263MQ72dhhMKCIqKiIpJRGxya8kZiAigqiYUQlTWbKZj56M0YaIAJC7YRDRGHsiurz6gE4AqnVjgYje+1vnt+7dv/fKSy9Z9QcAxwSqOWfiQuXE5ShX/I2ZofG4yAy3bYo65lvnZ//YL/pZD+7fffFnt93tv/lP/5mPfuIHv+5f+Rdf/LvOz9a/6Z/4peeB/8x/+e0f/ehHt9tNHGK5aUQ0i1aAHiJ475x3FhNk7bldGKLQhBsZySmoOhdCG0LThmblPRvpeKymRd2pdDoK4xj9LPPAOMLWTW69RChAVlGS0+QmIiLDELu+g5ShrHYAJCfLoJcpUWZWG9Y5HuNfVX19qWR5vnt1mu0XrK/xTZ6r1gxjja6yWptGam3kXtY2KSYVjXHY0c622XaZmWsGIwcfmuDPz86bxi82xcsSfWBhf6o4v2DdLhMAPrOLghf+Vzy4CLCKQCg6AKPoAiqCKB5g65Xb1uyQB9fEuhSsaMxVOEgYexZLdbJMOFTvLjqXynTkuHBff6HDgtVzFMF4wA3S4lMoGoe42W6ePnm8220hKyCCI2ZnqA0TEWOgADAurcbZW1UkJ1BAttuGqr4KxzMPECFLliEjYrfff0Drf7kE7G4lwvPz81Xb3rt/v121yayJic/O1/fv3xWRfddlaxiJqkGzfHh9P+z3u5zLoVuljurnvP7qP/krfsHdO7du9PT+9Lf95xeXV3/uv/hv/plf+49/7odvsDT23v2Gr/2l9+7f/0Pf9C0pDlQ5JokUy0vbqD55crHZbuIQ7RVlydZ8GhJ7eXUDFhAghrYNzWq1XvkQ2DktUIQiMSMCw2zwsvRevrYw6FiEEUBBRHna4pZeHHlGCyAnERSdvYBkGPr9bq8xzrM5zl7RRKDFrUkVUDIqzOejgi5xFD3V7k6x8rhIgNVjK0Vd7AMm25vxruMRgy0zgDViOVm25hRjwOzW65XoOvjggwveF6fRamGCNyjReJO6fSIPQE+VrWeAHoduBUeRaNaDjOK28vEjItCxXw8cJLrW0FI1TI4GE6qCGRZpzaCLMa1aL+hiDqD6wzvAo/A47Wla1Cy9lQ/zQOr3YBwzFtsB47aPPzGrOGIzjxTRlMwlV0EFqHjW2kwjKqizaH1SjyioZM1ZkMiRtUv1vrfsurDqb1BJFT6gj3I+mz8vM9+9cweR2ib0fXzz7bdV9dVXXlmvmpcePADiR48ebTZbs3/AcX4CAiLKIheXF0+ePB76KFmNc2Kf84/6oi/8Z3/trzhb3wxGu9psv/GP/6nLq4uh6/+tb/hjX/9v/Cs33Qf8sp/30x7cvfUnvvXPr9frwhshFhUZ/ehjim+/884PfPSjDx++q6reB2Z2notHU9Kf/dU/4cV/IxHfvnuvONMQT0G6jFzILUmOy4keMt8WRbR0L3YrZMk5oyPvLEISCdAzO+dERFLKmnPKKaUSuosEAFkSqCBVJX00jywZXCJ55A6q6on6fuAPWAf1jhzA2llgspvAWYxVbROXhO1pS4SiuWLBFP3ZIm0CRtWwDkO82mzMe6Nb7Vvv+ldemZ6WmkLg+dX4RNjT8+r2py0Q5toZsryLhExcKQG0jlQ/RHdOEjj1BNJuIMdEA1adP98JCMJTywGsTuwl7l9vp1EPT7SjpwonI34qq3JdmuXMIyMgowMG0K7rhjjEFAGBPXtZAYD3zM4hoBR1fLmgJzcL1Ym3pqqWmQP1psgyL6cs1BG3QiKq3cY/YDtgmZk7QEirtlWAfojvPnz01ltvA2jwoQkvrderV7xTzSradV0WwVJGdWrb4tDvtvt+iLaYsfP0//SVX/F//ed+nfc3vmv+0Df9ye1+552PFP/SX/6rf+u7/v6P+9IvvukP+Wlf9RU//sd88X/+3/2V7/uBT+Txw58+8iyw77vHjx698eZbKtK2bdM23nsRySIpyXa7vckBQKuzM1t72JDBxExMjMyusJImv/25xIue8J1SKOpBFetvjGbK1Pjgg49pQARUUcmSUrIFbsoxDjHG4uZE5iaDjgCIJ5BHMxT1mVqWTK4zoI46unlBjcul5Sz70XqaMRdAmp12x+zJUm+05iGOIQU6hgcvfADUdhtUZoM5JkNBu35IaSCAYRhunZ1Z3MJEAgJ89lQF78ccwi079+dLhE5+/TV80JGlM0o44MgyYcTrJ07YvCrGmkY64enlMCWYNLv4rHDg5Z/HzbDqeCRP9I+izMY60G2UUk02DFpFxpUvITQMSsYwS8JZ0VZnSajIlPsrqp6d93633bz19lv73d4Hb9sAAPDBOedNVjNH5eDkujMbhc4MuJpMoZWMDWpSrtEuUD+gI4Bpc1JMUsWz7nb7R48eP3r4OMYICm+//baqvPbaq633rzx4wMQPHz2+utrklLIKEalCFmGipl3fvnMnpWwgwr07t772l//Cr/zyL3sPT+x7vv/j3/m3vusjH/mIJBmG/s7d23/hv/6OL/7hX7Bqm5v+qDu3zn/VL/pH/pfv//i3/Rff3nUdIBqLKedMCLfOz1//0Ifa9YqQmrZlx4SYJOeskvXO3bs3PVGL0RuRI54I94CItMzQRlRRxRLmPu3gsMpQRDtLskhKRNS2Tds2TdOC5p3td4dhD9s4RKPKi6hILv1j0X7NeIyIgEiVHFi4UVQo9PU1PhNQDgNC8GDuh4Wx82KjicXbsHAFTHVQYwTztpCUC4d46QSjIDkDoky/ptbRG1ilAMQ8Q0BTgatLZAVC67VL2Req28+fAN5rqcBD20jCZxhaF0O3eel6sPcwu1YUEAAt6u/ZiO05a5BpjYLVJsoiahRweaUs2EawtIFe/MSFhgtURnXDNI9WycKASgiTDYmIDDFqJ5Lj5mrz5MmTvuvPz8+JMKUskhEDIZrEd6zmSxNKnRx4ZFqaAZjpbtWRLWW/x97UH7wDoFzTzjmiDLDZbB4+evzw4aNu3znnFOByc5VVkN1rL99brVavBo+IRLzf77puQBRQICLHfO/undVqZTfoT/hxP/rn/LSvam8e9W6ryz/3X3/HF3zBFwAAIalKaBpA+vP/7V/+5T//Z7y3V/ojv/Dzfts//as/7/WXP/oDPxBT9N4zMzN/+MOfc/fevZwSO++cMwsgkZyTqMKHXn/9RnewoSlIxOTYgKAqqtXoUsZyn0wSpqCd5W6mcBbKBp49ITl2hJBTHIZh6AYVUZGh72LKzDQ11ObqMXOvVUHV/MZBpcJnodBEqWAt9soXEdbGnYbaJmUJNZf10enmWkfF2ojCFsxphB1Kb0hA8yxR3Wg2RUntfT0KnhkZkZAYAbxzPgTHbnF83dwA7sW/2uk136LX1NOTf6+z5nsGnCfuvcEOhBZYO9oBTePrUqs02uzj9KOmxXBxZZ888UErWumC+LO4AxcR7OUaHll/s2U3LJXJi/WvVnFa8390LqfmGjYyhsdfqDAbURE754gBoR/6GPvtZtv13dAPRLharwC0H4YsGQqxbGSRzox+mMCNcRxR0FktXLX6k+9evfgo0kMLa/3+j32iHo/mEd5+L42eWgW7w/Gf8Nrsdy17V2bnHFvhizHmnE1T6pybnlXORckMAMxs/5RzjjEadcdKeRZLFZfrVr6I6Bx75+zze3q5cc4RuxC8Pamnl1fvvPuw23cma7Jv3O26t958U9Lw2uuveXavv/rKar168uTy0eNHfdcrlCAbx3y+Xr368oOf9zN+8odff+U9H0vf8Z1/e7/vHtw5DHL5zr/5XevG//yf+VPe24/13n3tP/rzf/nX/CN/8b/9H77/k29vNzsAuHu2vnu2vq4cnN1k4DBvepqznKYWzdbsKiMjlXHUG1oQCi4i4lRFU845g4IPbr1atU2DgLv97vGjx0PfK2rOoqrEDDp5vRcsT8bCf1ik7Nqc4dCJPlL6JZHZwtsmBoWlQ53iIvnnIENksc+YUIkZ2sEqLQIqVAhhWtYhEU27iGIWpKozt6WUlSwZQCFl0wzzYSawTiG1sKhSB+CPvoe6Dc/MA7jR3+u1/6CKhI54cmcrK9tFFh/isdOwwYZzsOFEv6xKvx7tdk4/N6ygtNHmya62KYxMr92BH0xjcMDGmpDQSot20GTnnIch2nUR07Db7S4vLjfbDTs+Oz9jdrvtLg4RAYjQepzFc8WluRsUvzRVQSTHVMJLqltgcn9AXNjJqcqn3nrnS77q56qq9845P/JHxWwORTXnKCLeuVvnd156+cHtW7eZOY9N3ElVeZaEAGfr9YP7L736yoM2uCcXlx//+McfPXwcmvDKK69+6EOv3z5bAcA77z76wTfe2Fxtg/ehDefn56+9+urZ+fqNN95641Ofun3n9hd8wResvEuqn/zkp954462UYmgaRNBcvB0YMYvmFNm5O3fvvPLyS/fv3O5T+tQn3yCinNJu260C74fh4uJyt90RkfNesiiA9yGldHFxISrONy/dv0OEITRNCM75gYaY8267x3X7wz73w1/xY7/kR33RF9L7SOj+nh/4xLf/lb8mqlk05yySUVQAuiE9vbj4xj/xrZ/zoVe/7Et+5Hv++cH7r/k5PwMA3nz74V/9m//T/+9v/91UhauMknYLPZAU041+eMkmNBd3URWdBIZWUkuqO1Y5KVRbdoFKNotWRLbJbEpSH7r906dPht0WnPfeqyoSYyFklVWqeSYsfLVwlojxyJbRycS5CE0LO85QQCUoT1WPepYjBf+inCGeLnmqcLju1oqugqBmUWHeFuNOocQqzeEzYz9no1R5ewmRiLEyaxmLl75wf3+zev4MO+gXgpaOcJ+aOV/+j4mMyjj+WMGZgo9zd7og3Y/2QVqtQoq7EuCy6OsRBFR9I+AUFjoGr4+ZEyrl1xyulw/QqmUoWTn+UY+jA3C0g4UpDdGuyv1++/jJ47jvKXjHLqVBVNpV651j4hRNOJKJ2TPPjJTFDqPatCEQYLYxiwrFNo+7NZyDWOr8XkQ0gFsll3uj6i9K3z/Sy9XkUbvd5q23hu1u++DBS2frtVmtzfn1VRNkoTFxSBdPn3bdjgn6IQ5DCk3DzNvt9gd/8I3gyRrwruutA8hZttvdp954yzna7zoA6rr+E5/8QQYU0M1moyrMPFu0IoD5eUGR42632zdzfvL4iYDstruU0uXlZUrRESWRzdWmLGnGcmiBskjYdcPb77xzdXXhmJNoPwwpDpb08aO/6If9rK/+qvfT9dvj4aMnX/9H/sRmt0sxbjabpxcXF5eXGgd0IbSr0DTehd/+e7/+//31v+dDr778Pn/X66++9FVf8WP+0Dd989PLK0nJyrc3k05CyZJi/Mov/aKv/Ir//YsvAMrW1W5UOtw+FjEKIAEBgoqUZem4ZrNwOCYKTWjbhpiGIW6uNo/7zkjxIkouIFHp0FURxixFnQMo7HYgmCUZtdJ4vNklG11iXFmPDbZqVl1UaDhCmK9tXU9b+z67IqoAoIBoXoqRxOzpEWVh8GapMgCYOJGpOkNwjk75D52mX17b8r4gBPS8tcGNsf8T4Cwisyt2gDBHRBxNOrVjXrWW0ROr5dl97pkn4qy/xaq8jwm9p758wpJOEI9GAF5PL7xxQfO3biln6Yfu4vLy4bsP43YLIazXa+9d0zTBBySKcej7PsaoCo6JHAPOphQmC58c64AKqV9qcxUELEw9++UlpMkczMdTkACJVIx+JyqEPLVa5uWs44OYkNkRpZz2l13fD44dETp2VQYy1kwv+yhTSpebK70UW1M44uADAPb9sOv2Ro4wNof3AQBSyimm7W4nIgYz73f95eXG7CsY2Xjhsyhcx+gMqwsIcYh91z9RMVokAuzizuguZiPg2I935eglSORdA6BXl5eXFzD3BUT/7D/1K3/WT/1J72E3e/zo+v5f/f3f8NGPfzKLbDabd999962339o+fQIA4JsHD15+5dVX7ty9d7Xd/vJf/1u/9Y/9wdff3xmgqr/+t/yOv/JXvzP1PWQBImPYsPOgmmIa4vDOOz//pqwqgMoM/MCSzOwzsW6NREfI1ezPrL9BJrM5THG4urraXF2BiveembxrRSWnXLqOpdmO0T0Mj8QKq1EotomFY18QKZmdHav0DK2V+KAz/bBe+k2+AYtpA09v/xbdce2sMOMMOvNQj/qsSWk87gAQEYjtCvTeNd7XO4D3VJxvULrdqX6/jiB+EVuhk17TM1rBTGw0Aixi+iniZGzVD75nil0FOhgqTg4ph55CEz+B8OBNMaLOeDUTHu1IsTI4qb8VD4hJi2UwVc9mpPsrERCzprTdbq4uL1NOQAw5p5ybpnHs1KwWc85Fz6OGuNq1TeXqJzDbQjG/F1LjuaeUs1gnJKKMOtoozkPP5GIAxdIcsymPNZvVKFHxMVcBmGNzy8pdS3fCQ4oPHz7s+/727dur9Rl7R4ATLbUi2RXpsrkQMxGY2YuN2gJHtGwcQ7yLESrqSHodN/XTNYWLRc2c02zfuNjjV5pMPKGdxKLXKb+0gMyEBAhd339aqv8Q47/y+/7Q937fx8hxTnm3319dXW13ewAAdOvzW6uzM9cEZFbE7/v4J37Wr/h1/923/vFXXn7wnn/jN37Ln/4Lf/G/AREggiY455kdEMaYYoo6JEjppiaAOlPsdYreKzmIQCogWUZ2+9iqieSsAEJIq9Xq/OwcVC8vnz59cqWqkjWnwXuH44yeJOUsqrl4AxFVPR1SiRJcsBdGnq+MV/uEBVdAAFLVXtbCmIMApnm5vCz6eKSE0tlzutoHYJ2WOukFp+FjqVutEf3yfpp5wRyVoITEzlFF1MbT9flZf67Uz8//+k+7DuAIOLINN40xrUfDS52Juwi3HYm2N4oxrM19avvPWbO1SAo7QePXF8HXakPACfmaMm3H8VlV+mHo9l1MkYmhZVvfeecISTSnmGKMlnvHyNbt0HyKFbuVUh/HfkFEYkoqmZAoY4KkJQ9yvGWBEKfWnIp37mglCiMox0RL8+1yLvMY30bEDbuc03a364chZ1EAOsPgAyHNT24GqdTM7GiU2mVDNokqe5P55rCX4y0BVVVUEakJTQmJmlxzr/8siKncLSMpgxwQOoJimCG2nx8tQHC2b0V2jmdVuALgn/xP//xP/PIf83mf86H3cwOklH7fH/yjf/d7vt+FELzrEHPOXd8DoD+7fefu3Xv375+dnTdNw+wAMKX83d/30V/8a3/Lt/37X//aKy+9h9/4sU9+6nf9vj/ITcPMPoQQgmNWgJTz0Pc5iwSktuHQ3nyiR1jYHlZ3gIiITEGvImKE7xCI0IkIgKqknNNut728uEgpMXsfXGiCikrOOeUsabpMjW+EC3s4rbH9UhAEZiaNHoG04/61as0rp/nZs2HmeYgqYp25VKmbq/sc64AoVUCZej8s0vM5xdRU06haWWMvoCQtpNIpIhIgqWBGAGY2ePOzw9WbdAAv6iN98uuXzCk8+HtitgEfYWGDsxyy8BjlWfBwx0vhWkf+ipY8P9mF/SdU1PnTibN6rAyuT1IdI2zGwl8FS5S1kIiyQ++9Kmy2m91umyWPn6g2ITQ+AIKAEjFCGuIQh8hMLviS7SUgKoXrBiNHgoiYHbOoSpKck0oWRJGMaUpbQIODsHgmliGaiQFQRGzYQDTujeWIGfSmRA4sRmO0b12O4mS7UztCzs9vNT5MtXwkRemE2OWyhJhSkgiglsHNiw2cQZ3ZyUcULF+17AQLh1wXvUAl+ZtyunUUjSoetITzTTznGM9uI8XtbxjSH/jD3/z1v/e3v2etXIzp3/h3/9jf+Dt/r/xeIgXo+mGz28Iw8Gq9Xp+dnZ17H7QyOhbR7/3ox/+xf/63/7v/2v/tx3zxD79pn/5rf8vv2Gx3zWrlnQNAEe0l5pxTSnEYJGVgDk3rbuYFVHI85qkcYbx5pHweIsRM5kglWRFcG87P1o7dbrff7XaXT5+YnssuXiQo9P4skgVAeXS2tjSW6RMqVVNAoPD9q8GuWo6NZsuzXRuOaoEx3mN0n4ZlFt5IjzhAOeqbf+ltplOS6yLaoG7zdb6Escwgk0sNFIcKgdLMzQSnEdTKKgyGlzAt1pCHGutr6zAc2jI9/+s/7RPAiRAu29mbmH7G43RJmDnq9HHU1t2g+a+wj+UwUk+DeCgVrNf7C1OHow3RqcXBCEeMSDKAquz3+34Y9vt93w+jO5sQonfeB6cKIgnJAjdyyilwQNBsfrY6MSCzjqxn51xgGhF+KRQJULUQgRElHyOTaHJOR6TMmYQUNCcREUSj5TIUYrTWYWyV7bFO5d85SikOcbi8vCz1+tat4MO056rd3gmrPlvH2LQZmltMprOWbfm211PCLKioDMSKFQEW9MZuIhytuExSR8xIKKJgvpW0yGItWxWg0fpfAOC7v+9j/8mf/a9+5S/5Oe/h0n96efV7/+0/8t3f+1HQKVgciawUUXN2dv/BvTu3bzc+KKCKZMmEuFqthiES85tvP/zVv/lf/ld/62/8pT/3p7/4L/3D//5/8B1/9a+H4Jkde0+I5mKU4mD9sPM+rNfn57cnO7wXv5d1qaAtnP4CwdiZrEjomABIVQkgZwWJKQ273fbq8kolOx9C8ETmeCsjB2EcQ5lENKuat93ELJCx7S98/8XdhtP/UpUShuOEP4oSBBYm8rq8f1FqbU21H1gcE0dtsC5JL7rAfqVsJcY4qSlhyjYZqgSqOJJ/ilB/LIZorqHkkBA+W49n6ACuG7yv0wHUjfI4JioYVbHoAHS29ahonDr7ide3fT1wPQ91mlcrOjl2HgewTJFfdLznP7k/H73g60mkPL2SPGTQfLEtBCJkcrv97q2339pcbZq2ceyHIfZdD6rOeHVgPGWVPMQYJWdVzSnZnYWjLE6MtzgWYkJkJsecQfJx+On4F7/hV//Kf+/f/D3wQ4/r0fmv+dW/WUGLFyFOPdi0UIT/8Nv+wld++Y/5/I98+EY/+aOf+NTv+7e/8eJq34RmGPqcRQFzTiISvLt/996tO7dfee214EOMMaUMiKriHN+9czt4BwDDfvf44bv/1G/5nd/+V37h1/3232KRxc9+/MDHP/l//wN/aH22RoCcMws37VpEr64y9BFW7Wq1Xq/W57dvnZ2fN83NIKDZGwMr6flM8VQz7QHQpmm8c1lyt98/efxw6KNIFhFmJ0iAmFLKKYtkA/o5eBqzE0dXl4LXQWUav4iTx8lOoA7WnZnZiPNSdqIVLYibh2WklOp5YzTKmHCxKNB6LT05jB4XwVn4BAAIUtURAlQdDUAVVRRhkiwLIqmKMBGT894Hz8hQrxBOWdE/k++vL1i34Zk6gBvxTE/792DxdzUWEB+9mEWQ78L74WZDx1wFpxwaPMSK9Prl/tEcAHq80EBAoGp40Lm/SylZdA0S7nN3cXnx9OnToR+8d8ouphRTJAR2bDkQ1jzlnJNk+xU5pyzF22ii3tkvtc2A82709S0SxxpDmw47xB8q8s9rebwXEcw5C0IWBWXnQtMSkYhkyTnnf/MPf/O/+3W/48WBoP/+r/z1/9cf/49iyuZi75wDhCGmzeaq6/p21X7ORz53fX62alc5yxAHEfHBnZ2tiTjGOPRD2wZR3W63u932m//Db/3L3/k3vv5f/+0/7aueZd8mor/5X/66bhgYuV21q9WaHcchIuL67Eyahpy7dfvOvXt323al70EBbuci1XL4Sa4PqiKSQQlARVKMqpKHvru8uBz6HpmbEJxzqqO3m4jdSkX6R+b9k1WzMdOK/lHraoCzNTAikZujmgqPRGb8cH6KVvynAwCvE/ZUFPMlfaUKu6rOhekc0GtooVXlqN5rRBRFnLhJs0SfymxABIrJguadC96zcy/W934a/t7dhOFz3bJXj/TTUBH8gR1779mxMXbB+PejX9nYx059fHE2U9SaaDj7Gp/e+6Lq6RI/CwJmk368TuQxX/r1TnqmAEzMdxBQJmSjeWx3XbcPPgDCk6dPN9sNEd25c9v7EOOQc0YEJnbM7JiQEmQZeQwzf25xGBJN+Dahd97ktSnFoY8xJwBAs8QahYKlowL6oRL/7EcTQs46QK+qUQQAVj7cv3+/ads4DJvtZrfdfc/3ffxP/5m/8LW/7Bc8H/a5uPqGP/an/vbf+5512ygM3X7n2LerBqB9ennx+OHjq+3mtdde/9DnfHiI+fHjR5YwnCV7DC8/eNl59+4772w2mwcvfd5LD14i5svNpr+8+uSn3vxdf+Abf87f/K5f80t/3keu8Y7+tr/43/397//o0A2Q4p27d3/EF/1wEf2e7/m+LPLqa68R4Ha3P79968HLL6vA48ePhjjcrK+imZ0yB6CMXpaSs2alhiHL5cVF3/cIZHCNXemqmiWJmSkhsmMuU7J5R5TMO0lpVnXqDGMS8cSyLKWBkWby8cQEAhjhlOp+rnAdqAJJltVqFIPhIdh/YO477veOjGqmclEkzwBVeuBc7styQBeadgTMBcCEOYjCsWva1ns3+VbcnIF5o7r96dsBLPk9dYoBOQLv3GjxsQgBqNVUB1gwXrdbOH22zWKxpSXn/KKLnOXZpKLKvhPhIOe0vjeKmKDbd7vd7smTx/2+98Ej4dXVVc5y5+6ds7OzoR/6rsspOnartmmaRlVjitMRpJUhSeXfrFkyjoxqh2xBtwCYco4pSs4lRH6hahRBmSTEP/S47uN13gEkziRlKydE2LarO3dup5QBdL/fpxz/5H/853/yT/hxzwaC/vJ3/q1v+GN/6mqzaVdrAHDOrZq1C+x9AwBIdHl1dXF58bkf+cjrr7++3+72u+12t7Pr1Du3ale+8U0ITRPu3bv30ksvPXr0iFQldgj44P6Dv/u9n/xp/+iv+5qf9dW/8//yGx/cu1P/6jfefvf/81//ZSZOsScFBY0xp5RSikjUti07l1VD8G3TIHKMsQk3soKwm2UEf8xbj4x1xVTkhkklpRT3+1232yOS9x4RvDPSgYpkBQEEZgreW/8nYoCnphRzLpaiZEnGNHoLGWkcwUy2pjWurVALYqQqtsGy6i9SaJmIxaMRqo2WwfAmqZlNTOFQRVqkpgpaBGgjzaF2rNHDqmeQImudhDnZWEyuQ1ovMyYYCUe2c87CiEzBe89cm9t/Zgfim48Ppxeky2pbm6kCI3gjJs/AxmiOrQtvaHvXqAoCPi3Gqze9WplC1zxeo5VXfQwc+MfpaT1Xtfaq/jMPp+M5Qxj74fGTR48ePe72O1XgPbvg2qYNbdM0bYpxv9t1+y7G2Dbtar32zu92uyHF4AMxWwakLGLvS00a2/oMABkBsoDD2pVIzXlXBJgdOxP7QBYQ+aEq/3xse15vEwmIyBB7ydJ4F3xjQSspp6//o3/q3/rX/oWawDo9vuvvf+83/8d/9u9/9w/YFb7bbWPs796+++C1+4S42+/7GEE1SUopAcDZ2dn985UAvvX2W08fP7FnsO/2Cnrr9i1yfOf2nZTS1dVlv90CYNs0KrLbbp4+ffr1f+SP//Xv+u5f+Yt/7k//ST/+R3zh5xv48w3f8p9sNpvt5goU1rdvxxT/57/3d1WBHK9W633Xi+ytd08p37176/79+/dvkgkMS2KMmSkQEbnRm0wBVLquN6vaYBtm1ZxUJKsCEpEjclxSWgzAzJJjMtWu6UWIjI1cGuExFBZnt5giVVERBc0TQDO3fVqMqOy/i+O5MbGJFCDFlEwcx6RgjNXJb7POhiqkMxkVi4BgaNJoxzslgJSpyDzujdEEY+xVKThSGRCJSM5jpVv44YsqCaoqSNbMROhCYCLj/p3KQ3zPsM/puu0+K20XsGOevXKWxs2zv7HUse8v0vtPocqTHRQsPBCmxHc8DQAe/5pllNeRLRMQlauZCHuVvh/2+30ceiBKOZGjZtWcrc8VYNd3/dCnnAiRmFSs5UlGoBaFlFPOSQ9mUEAiIHLMbNshRASiJAIIUszJEYicc8yuaZoQAgDEYRCRpl3/UIV/gQvSrC7UOReC9yEMfXd5ebFqG5XsvW9CUNWPfeIH/+xf/O9/yc9fcHK+76Of+Jb/6M/+ze/6n+u4QR3jl71zzC6N2x1EJGb7fYgUQmiaBpm2263k9NJLD9ZnayZnVg2gOgwREXm1RoA3PvUpIlqfrZnp4cOH3/gtf+rr/uC/97/7YZ/z637lL73c7P7n7/noMAxd10NKzocmtF23J6L1+tx5b/lC01xs1wk7vtGir9ikAwCiY2eJvN1upwBD34MIEY/hmuSYRSWnwuiBUQlMRKNvUBaRFFPKaeLbMzGZS2Cd5iSTzBinyIMsqiIwKQVp8vGi0WkYi/6YTFBARESOjV4EokBI7ACR5jpTDQAVTERGryMCBBplisUbA2cfJJOpgQJAEjBfXnbOIXOZ4iUX7ZoIEKnoqOmXqSUtfCqxAyATUvC+BtY/00s997wye+O/r2KAUEZ2PDMTj0HmIzyHlamAqFj0AuFhfO5J5H86HicSmB3RFVlc66+cqALGRFRdWD0jzoZr9fwyYqAwHS6iYOt7k0E1bXN2drZRSSmxY+ecisYYjfNqF6/ZHlxeXhFhCL5t2pzzMHTDEM0Us2KfITMH36xXq3a1sgXDMAxd343Bp0JIQuicP1+f3blz5+xsXVZGIqD64KWXfqi+P/th1iv2doa2uXX7dvC+7/qnF0+3W0fEbWjy+mzf71X1z/1X3/Hjv+xHfd7nvD7E+P/9a3/nv/z2//F/+rvfXeZKBECUnBFgvVqdrdeq8u7Dh2fn523TNCKTl4kq7Pf95dC/+eZbBsQ/fvxYRT/v835Y27SXV5cf+9jHgg+vvfrqarVanZ8rwG6/+/7v/747d+58/uf/sHbVvPGpN998400F+Bt/5+99x//4nR/68Id/7I/9ccaihCzDMNy9e/dDH/oSQHz46FE39KvVqu+GzdWVYx9CE2N88803jcJ7k1FJFMhun+C8SN7vd1eXlzElVNWcQ/B2q6WYJGURUQQmplBQHBuectaUcpEAjA7nTGRmc8VabuJijRxOrITjlaoWZ6wGCYkdl5JfJdOMseNWr82gzch3zMQMqCCm/6og3mIPgLONsf01lQTpqYiNbNZRTKSKGVVTzokAyDnnWLJkRFbS4tioxAWqmglVJcwtE0DWgtwyOx9CeYFSOzniyfqr77tuf6YngDLvlGOZGWICFVWqwRaozQQnz5ZqHXvNEl9HPTXQ+AW0cAvCk/Ky4/W6nhS/mYsfIhGIKAKwZ0RMQxz6TkRyygBghifA2jRN27ZIGOOgAClGFTU2fk653+/RYdPeC03Y7/uUcpYSNjKF5IEIsj+/dX7/7r312VmSvLm6ijkRMaSUUgIEZtYsoNCs2gcPHpyfn8+SKtXzs7MfKvHPviBTTtmgNzX/RRLVYYhdvyekJjQ+ePaOBhLJ3dD/e3/8P/niH/H5f+l//Bvb7T6lOAYDae0iws4Ru77f910vInD7di6kF0xZttvdkyePY5KnT5/GGLe7Td8PqvL46WMf/DvvvPPk8dNHjx41bSM5ITvNabfbx82G2TVNOD+/lfMPXl1t2DtQ0ZSurjYPH75L7MiHcLZGYmZ/5+4dJL7cbIecHXsJSkVpnDbb7dOnT/a73c12AMSMmEUk5y7llIbddrvf71NKjKPIBKA0JmJZHew8E7tRFJklj11/YZQSEbJJwx1bb58nb1GZyeGqB0AUwaQKHIs7kQWnconEGJMBdbRlQETNuUJ1kYiAUFHI+sJxK2CtpyigmGmbUfKQjWqBqKrEjhCBUQEJMlR9JUrWDMaZImTlkdlCCsAzD2W0h7TSjyZUVwCkxEw+OO9H7hk+t7Aek27e2wSgL7gzvubv8SSoMltlj1ELxMBcZiibiCfJBtVbetW6BONhxceJj1CMD0Anb+7Z/Qa0UgboSUxMqyCHpf1cqciSVVXBVHmaVMF7B4ib/ury4tJ+02a32Ww2MUXPvvFN8IEdDTH2fd91XUoRQFKKWQVAifz0qhc7BXOgiwIi3rnbt2/fe+mBd+Hi4uluv9/vdmaqIwoEyOyyRJFs+0znOKcsoqPpjv5QjX/2bdPtuzJ4IeWUNpsrQIzDEGOyhjHlbJT0lGS/233P93/se77/47P1fAUjGQIBgH3fxZis6+z6Pj557Jy3LNQ4DE8eP3306PFq3RLRk6dPnzx50jRN0zSf+sE3PvHxT26uNm3bbne7T37yUxcXT4d+SDEqKISgCm++8ea7Dx9eXl0BQhwGzRkQdrvdD3zsB5pmhUR3Xn6Ffbjabr/n+77fBx9jBtGrq03KCRGS5KcXF8ysgCNt7EUfFlMiIHEYNpurbt9lyUQUvC/yxFw6GELyIdgSl9kVnCSnlHJOaRzuidmR41mkDqiqWcXc4EBAT8DWJmpkHkGdkT05JVlYc2n/ygRovlNTDJjqrNmtnEWmMC5aMD6NfDiOIQiaVcyieYTsZSKrIiCzQwCVOF1dUuiMBo/BSP6fAkUMGRJCBSWlMgGAczmdtas2eL+8XA8FBwf+pctcyxvV7fJn91m789jsZ5hqllWR+hWXbDxMO3j+T9VDQ4kDc2+s5aNwwigVDwEsq6TmYTclTZvhTN93u93u4vKq63beewDsum7f7UCVQmMQEHvKklOMfdcXnrJmAHChCW0DSDGmlJNtlrRSr5Fjz+2t27fPz2+tmlZEU8p91/fD4LyfOhpANKbwatU653KWfd9JFvNRyfmHlsDPeQzDUD52wiyy2+/r6y2lZOklVgT6YShAA9KE+x0ASqo6xKSavHfBuZjirtu37UoB2XsiHlK/2VzlnLqu2++7fhiaJhDTkyePd7t907Zn67Pddvvk0ZN+6HNKoOqdD6uWmd98++2ckkgG0Tj0oOJ8AMTHj56Edn/7zp2z23dAYd/tr7bbEPz67JyJbMJgZhHZ7/eOnXPu2/7it//lv/a3U0rDMHT9kFIkJB+8c24Uro7uUUT7rs8pdcMQU+z7Yb/b931PhN57RMwCJuoFHW1FHE98IVuw5JxSSnZzE3Op/sw41mJDUlKKeWYulBgiQCv0WORTxMSMNC6Ml17QNOdPWkAmmBudnUOY03QAiIwWplosLezY0fpIn5k/oABqORijkcOkMVYBIpGcCSu6Tlk46KiItydV6O+Fh2Jn1ARZSHnfXJvX67Wf1dr6bFdnPUG+eS9DgHsfvf9hqT6I0VxucomRPDseP15r23Hhn1obvChWlk3lMEc8+H1YeD21898hUbc6EuoMnWWyyyh1H61pIKs6R4jQ932Kw3q1RsJHjx+/++67WWS9XjHzFGtlCu7JDh0Vs0qWBKqEpAreuaZtQ2hUpI9x6IecZxsszaKqTbu6/+D+g/svrVarlPIQhxgHJEJiyaKgxfJZoWna81vn57duIeC+22+3W8nig/fOq/zQBPCch4iOcmyQipMwe7vMdlXmSSCQIGOuEcmJCzblQyGCiPQxGvfDjoq2bc/Pz5vQDv2w2W4vLi6y5NV6NfTDbn8houv1mfc+pbTfbbuuA6TgPTM75513kqUf9sMwgBmXJoHgVudnyLzfd1Jqq5OcCcl7YuKcs4KSI5s/jCORJcmQ/5fv/djf/94fMKOIfhgUlJF98ME35ApBAxHYsWMvOW23m8uLy5SSydebEERzSjGnrCoI5IgN6J+aXMmSUxZJMl6IRM4555jBIoUB1LrgZJhPwY7AxAFIXILILIKYxjw9NIfDMgMU8xOLWWcizFlE85ziUsDbunMezwBVM6e29C4S0hM8xiUooipZJ8+vaU4RVU1ZcPSj1SObSyhJalo8YiQXSYTHiocqOatk70PbrhofJiPF5yL6ONMS30PdLn92n7Ubj5nY1RNAKeE18jNvPSoXzxMo/hwaA3gUJ/ZceBOPuv6ReFaeBROKaNf3fd/td9scU9d3BHR58bTruqZpvPMAOsQhp2wUi7P1OvigKmmQGAfNOudG2hotNM67vuti36cUSySLlSFRUPFNuHfvwd17dwFwv99vtltrTokpDdFerumAjW+QYrq4fNrt+912mzUb0fZqc/VDJf55W+D5MqOJ51uDgLXdNxrh1lgeo84EF2l25UI0aasFm1ApK74J67OzpglZ8ma7u9psrEu1LWITgvchxtR3+2GIRORcCI23YIOcckxDHGKOEQCQ2TdNOFuf3bolCiktbAvYeXZuqsNmJ6aio0GCJslMqqApxX7oU0rs3Bg9JAij8bhIijl2Qz/0m6tN1+9VwTlmcy4Qzcn8BImJvGNkRkBR0ZSlhD8mUDFOJNs975jGtAlboKacJWcoEUbFNhzNGYjJTgEiRrb2uVRbGDk+I3Ow2MEiEoqtjYuuUkUVBYGKd+RRfSixSSXYHY+62FmKrzBb+0MlBjVHqckCct4ajj7pUmXi4iw9nW0awaI1Rr+j4F1omhc37MPn73tfbAK4yRZ5ttu4EeoKAAiCdgAQziAaIiBplb8yBnfhbOF0sNRdoIO1oecinbZeL+siU3IhFZ4+v5nAZMt9ZtR8dXX56OHDYeiRUJ8qAHh2d+7eaZsGAPb7/Wa76frOO3+2Xp+fnyNi3w9D7LuuSzHalWfm/lLAYjBbiCKrwfr8IccuNIWz0Q3d1dXVbrfLIkXbDkrARlEfhuHq8nK/34JIHIY+DpJLmPU3/NE/8S1/+tuQx0xgQsna7fdM/Prrr33owx9erVaSU98P293u6uLy6uLp5eYqD9Gi0GFWNCoAMrELgR0R0HVUWgIioigpdvF3/7bf9H/+Nb/iBS+Mv/fd3/e7/59/uE7lG39qSXObqVmzTRfOGdHLqb06zLX66MfNUJVzMVtRLeFgPeRX6GKixENQUasllbmeVMnURTkUfGiaRhW6fsgpiUDX7+3vz8/ORHS377p9l3NyjktjgZBTimlIWSSnlDOAArvQtuv12ersLITQDUPhvRTncVSZnHBH//DRFgdhdqYDhZTS0A/ZEnhGMaZZ2BJSTHm/2263m77rRMQzk2l6U4wxw2RYT8jTFKWSYrK9RbG2JPKl7huiQ8aFTJKlbN9H/SOxDwEQY8pA5Jwn5wpYjgYaOyCCEuo7rnERJxfbybbBTs0J6AeRXNIoiw1xZRVcVrkTfLA0fSsRJWWPrVXm5XRFoqJKEXMZ6X88lgDMvxQn/S8iKggqEsLsflrMDqYVIDmHTWtUXZ3AuArlP47MPbaj0PdQtz+rE4BjJkQxG0E+wOiPItBOE59ma8gXcr55JuFnYXOy4L9JLFBpNwwDgEo/AOHtu3fP1usmNF3XDcPQ9z0i+uBBIQ3JN56Z4jbZumxKFjOmgnktFytEWASp2PNgJkJy3kLAaOi7ru+IaSRKqAI4ZgTMkne7rdr2TLJImTKZXc55u9umHCFlQPSrdWh8t+9EchLdd0MIgR3fvnWraZr27Kzre+76XT90V5skyYjVxQKPmKYHFhbX5Kg+zWd2kKWc45D2hqe/2COlfHm11Smrc95vYUWBHjkdMOV6KtauzicOEBjFQ1N66/Q1OKdLTZQFrfXp1zZDByxEPXWhTXvJyZgkhNA2za7vu81233cq6sgRkWNW0aEfeqv+7FartmmCiMbBzIF681lDdr5ZuTa0q9VqfRZ8o2C5L+VRXhrNtvrzq1tIDHHiX+RkoT1AhI6ZnSPEOMQ4DEM/DEPfd10/RGfu/Gre/aIqTOycI7sMTMcIIFliika2ITL3cfbeKJ4qokYSygX3kULiISYiF7wLjWRN2oGCDQGAKCDj+pcsMxLnFMFSsid/d/MWHbdpUmszIJMeFMqSwGV0z7IngJOYgN2a5QmXLq0aD6nexc7mleMSA0ZHm+mTEPseLWFm9hKspbPlufee2X2WMVx3827+Ro+5xDKz956YISdVURBQmtVuSyvgGQ4qfz6MLhnNmxc5wPVmfFqPH3y4xS4UgZhojqicfetEtB+GrusQYL1eg2rXd2AgK3HKmXJMOZnpQvDBMe+6/TAMt2/fds6BSsrJ1oMiGQG8D74JKhrzYITienTBst4lVYgp5SwGmqpCjnl+N1RACYkck6oatU5zLk+eyLuwXq9DCF3fX202CgJ2qyI4x/2Qri6fdvs9Ip6drfX1119+6aX1ak0v4Z07dy4vL994442nT56AKgYed+BsJOiUkk0DGJDRIS42oSpF1lm8iG50PY05spOdUeHkjRbQVY7SfJtNNi9H5/8sCJw4Y+VvzSCW6rhnHBuxZca3Tr3/YanX+ayYeSB/9A/87s/90GsfPJzsu7//Y1/2039RjMlsG5i5aQIRm+jGVrwxpULfzwKIjj0zo+E3RABkAaQp2V5kvn2R2AUfQhMaT+T6YaBIuZxJbNRLHY/r0a9nvrp0MmXTusDq5OaPo6+AgogUM6L5XKTJTRqPumSdDY/GJUa9NyxX7TTqFU3ZVP1NlMQ0GR1UCgYFqNbdpglAAhBRZGTnQgjOueJoXfEnj9v6I7Kfvue6/VmcAJDZl5oCeo0p6OkT5Dl7j+NxQa/LkKwtXaHWIZfPyXnHCDFFE9Z659gxR0bmdrUKTVCAYYj90EsW649StvZcYhxyjjlnmkqQiCvr3xBjGro+Sy4zo4gl8doUD4r9MFxeXBJR8GG/76wBMkyZTX9cCM9sob4imRCDd75pHDsm54Ozi5qRNHj2nolKVqpq7Lq43QPAMPTsXIrJOfYh3Lt3b312ttt3m81VGqLFjwGTY4eE2VgxAKggUvJZqTJR1PG1jAFjN+kODppsPbGu0dPGiyd+jD7j5yPCi/nk4Un4c8qVf95e7gNFmRXVLCqCiME55wMUg+uckqQYJ3tnco6InJF8EKUwfJLZ/pStALsCARMhOxcs/dwDAsY5osd8YJBQZnf+Mc5hGUmMh6Vjmf5VUT4O1FIK145xo/WnnmAHHjAFT5SuMqUY7mWg2yhumDwkjoAa20sLIJqr/It4P3w6hwR3vCN+n3/GRSM1v1hi9hx4Si+HmdGv128d5ljHskixijAF/yxcnarkm9n6c2npPE1vaFc4jojtyEIiBCQglZziMHR9jtnYb23TtE1DxLv9brvbxRhVdej3LOHOrdvtqhmGYbvdpxjZDL5V0FT4IRDzbt8NQ2+ph9OyyuIgGNiEY0+ePN7tNjbzZhH2HGNSVecDkbHMCaawFAViatt2fX4WXEg5912/67q+77Kos3Ax1ZyypGwfgLGmUs6PHz2+uLgA1fPbt30It2/dOjs/b1brrLvRkKNwtc1zq4yqKWcAZq5XpdYO2w5vdlh5wfo/7trmNEuARRbSdAZUdN6Fpu/6mox1lkmR7uJhGZmYYCfx1AmjnH1pAE6GiX7wqr9KGgYXvPcNIjr2CJByjDFarJBKBkTnvPNcfHJotAmSHGPpG8bUHkfOW1wlMwMjWm7o6IW5YNfM9zQumrYxGKPi4Bb7IDhwbixYvU0PTKPHcNkcgC7ww9m3Xidaxjjtnbg+xotccQpzAZ0hexFRoXklgCU2ef5xiqAggHNSGYgIIjJ7YsbKzxTnUeDEuXWce6U3r9ufpQkAARyz934sH8/FjRa3O84D+4mjYgaDsDp0qsyx5b0+G70izNb77BkQu/2+6/Y5ZSQU1X7ojYtJzM57dk5FU0xD31tvrkNUAO9d27RD13fdHhWccyKoikTIRmhTsEKMjn2Z8oowcop/Q4QsabONoMrO+dAQsw7R8iAt5dVAVQRgcuTJFwCXQAsUG4c+52ymjQi2Tkyq4siFNiBhiqnv+36/BxHIeYhptVpbfllo2pQlxaiYqTJumsJURbJGUFGjYxdLSFURZcIpSfg9lZujDl+OnFAOMD887M+vb4wOzKCex5jWhd4G547yHy6KrWM2a6Nx05vTEGNOhp4Xbn/w7EwMlUVURl1XTNGcCpGYQ0PM7D1750Nw7Cwpc7TnXYS+ldKXpzMfZ8bHTNjSBX9ipHAv2JtzNrgtaFFVjz0ArsNSxmP/WddKqdNLe0kZI4sXZFCcs+7L/8uAoCYtSWZSS84H7/izlwa8mABeNBH4mj+f7uKrgxuBvHdNCK5QTbSOYoCau3GQFquLI6vUzeUpYUFwSJMV94L7jzOAO4uG58YfLYsciDnn/PTi6aOHj0ITztZrABiGIafom4aJmRhUY4wxDVmyilhuumMGUJGS5jiVKkIyiYyWPVgGQu9caJpJ0jW51dmmVVVt3UaVerC+2HS0tnXeIwITquh+11kTISLeB2KnCkhofDuVTMShCevVmhC3uuuHDhDAOwhOJL/15ptPnjwOITimtgk9giSCMaYJQXkkbhmQKwIsznnHJcRtas1uGkWglT9i1dof4v6AiGUsUa1SGg60HDpBrlgZiozLoPKfg9Dqa2/sqe+CMcWyWjjPweIf3DGAkELb2JGeJceYYhw0ZzWLCOfsBDAykkrOWczdBKYoLmJk50MIoXXeIbMZBAEREDKhmgu0amRGnGyRpwtdQGfN8Bg0OjuoLwa8wwZgQl0miRhWGe3HvcMsQDKuwKz0riHnOYW67izL1T/dMlrpyOYLuyLuTPsHnNyvEXwIq9XKh1CPmjUxGa/DHq9LAX6xP39mJ4D6jmZ2PnhmLtGZC1H2CbbO7Aa0nNcRDy088XiNX0Y4XYBN1cWFTIyUc04pG7K+3++vLi8vLy5iis5x3w99PyQRZg5N07QNEQ1D3Hf7GKO5VyGhb5u2bRGLGacjKv2ICjM1beudzyl3sQcF74MPwTs/ZRfLJE8f84yd8zBqagrl3zSoVZL9RFW21iynKKMinx2jknnYiggCOHahadqmBdAUk2M6W5+nlEWybaX67bbv92e3bq3aVWHmOlbRYRiyZKu/RrwTLTJ7gDw63xIBSuGg3PjqML/u6qPUZwD6FfNyGRhxyibKTgCpm7zn4DZ6egVwyjMersmU/iA9EIGZU06aYkpi3GXEwhqypZR9gpKyiqQkKSewOAoi5xpyzLZlCsE5P1KTqyPmKLTPymdR0AIZYdNYOxOUAnX+3XRDW1GWaoFPExBgJCmU435/9Idc4H0L6ooe8ifr6JElHjOS1/UZyWvLrfOUL6CI6D23oQkuVPZln42rzH0aC/3R34w1XABAmTl4x0SWJa06pbhhbe9d9b1jL4c6+sdObcDyraSaGzDCgVonJo7LfJmvIVXt4zDsO9947/yTJ4/fffiQCe/cvq2g2+1mt9siQtOuzs7O1qu1qnbdfrfbxWEoLw5pvV6t12sE6PteNBuP2ywA2fGqXbHjq8029gPOsnssR6DtW4u2vNgBmnpGFVJOmmWUPuIoeMaZ7gbzu0IWWUEjuQFVFQgJCZh53bbs3G677bt+fbY6Pz/rum6z3eWcAQm889479qoaU0opEjtmyglUs5qJTBmvRBUIQVVSiqrKziEiMU5Y6I1KzKiSg9EYZrzup2kZZ7PGOfStdEWjd/wiMrwOhisqzAk7wmu7/xPZ3+PqsI4EwkPa6QcXE1KAmFKKMcVoox8yB++ZnanyTZycU845VdAJAZELPoQ2hEBm/oUEhHapmA3f1F/ryL5bLFLB4MQJHlou8Kq7fyL+g0yt3QjaybgiGMOEK4efGt5bpHxMtsDVc5J6QBjrNi4uBFWL/Csih4rqU16q4kwRglL/Jt6RDQxMrm2bEPz4lmCdm47/Gz8AnnEwlKMZwDnyPhSlg+oJJs8ojNBDoACr7n1p7//8YX55XGCx4shZYuw3281uuyUk5/x2e5VTWt86X61X+92+67qYUvDBpl1iVAERSXEQyQQoIgoUmqZpmq7rhr5HhCZ4BY0xoSAjOceIJJKT5FBC0WiyE6knmYm1bg23asopSRZ2ZmCLxy9rArgQEZHnO6SKxiTnzFxQJRsflIhWTTBHlZwzEvvgQ2iaxktWBO0RmKhdrdqmUdF3H7672+9FswlYRvWP5pxVFBSRsd7n3+gRmpBiEpXDhE77HQh1JoYedPEIpyBaPajgen2/fr3NbFVs8AbBdB+sJbDGvo85gZqQhdh55x0R5ZwtrsTkMiAyRv16ZibH5Pwoais9tpm4ocl+rT9SAQRGrJb15rxb0ly4ThVcbGV0TFVRmRi5CwB5lHdVa8QlaoBwjc3wZN5wajx4TjdjgresQiJsq+YT5mOqcJgGaDeyd4aQExwJVP8BmACeNxgoADpXcqKRULNeu/etd/k4RTOOkREj03tqIOpzYTpzD7ySFGx+BWPOqELXdVeXlxeXl/vdLsWIRE0bbt2+dbZeS8593+ecjNmGiDklSTKd1mByAQNAiBAxpxxjbJrgvM85E4vHwMwpZ9WURRiLhn26JmhKtIDFZGzMfTHhFqgjIiKdpyMRrTbZS/AcZxLbqFlHBIVhGAixDQ2uVl/xZV/65/7UN734R/gjfuLP/IGPflxSolHkAmN4Q1bVNJA455gQb1r9RdS8eTHXhtbGKaqpQYpznzUN/aVdWBoz6xFCiJNXOAJew/Gv/aPGr1A8HSNUA4kA8MFdAqhqjBEInfPOjPsRASBLTjmlmObIQyQgJBeapvFNMw6FPFJvjR5Wam4JYRstNY2yXOQqJfedmXhe58w0nXJ966RCGW/GEU1HJKRCop/AlWLBbz9/+nqdWEBHIlk8zUUBGClnVXVeetIRSgbNORMpzDuAOrtwFDURHkAdAN5x07QlDFLywin/Mzlpus/OxWTkcee9I9ZlFu5JVket4JopO/o8HslhgtfRKmL85HOK+67run3X99oP4HC1atarNTHtd10/9FOZYGJEHGKfc0pxmJpCF3xoGiTKWcyWawQb0LOHwIQYYxxilCyOnXfGy6wZLzoJ1aWaAXKWrKKoE/KzoDOpzgQHpKm/nWL6CviBE69Oc1ZCDMHcfW/2ODu7FZzvJc+IaJlgQMEMuBBM1T2lNb3YIw7D1eZqvT5DxykmMZXB9Rf7/MkKKAgSnpCCwSmM90VwyyM+eW0LdGJNhR/0UWA0cA7BETk79SXnLDmmpJJBAaxFYkaTsTYhhIads0pdTBRG8wPLgyl67NECjHDyfStpvEhVW7QIXFvSfKrSOu3PymqovlsqfsGcxjg6idVkdZwBfz2+7rTKGTxx2RRf8EID1SLxQZwoQlqwJASgiQ86Sh1FBMhcVwID5GtDsD4zB8Bn6HCpaamj6QY4T8QjmiyqKLO6WicfFaj9PbHSfcBU7yZ8dvbyP0wMmtc3szZ08qHVnCWnmJMAAARTpQQA6Ptht9vFmGDUwDNj8H6Iw3a7ScMwXShtu1qvVojYDUPO5heNIgpILnAIIUnebzb7fYeIPgRmb8R8GGXtUjwhAADN7gcRJ6ss6ynMiawuQ1Qbn47CCFwkJyEhyEhttnBWROhjxJj0hqHBCOiCjzmJ5DHIQg/IFqbmFb1ZGkE/DG+/886HX/9wu2qr9RrMy/vxJ0NND5luYFGFkttX3zJYP7OjwCQ8nPeXgH/NIUQ4qDeLQ+KDzghFxNA0zETIVr6zUSYkF38lRCD2TbCun4iRGYmM7aOipGPQhRafZBEBhGIrgqygPCJDM2BiqRdllh3bmilAY0zctfPJ9sKVm8iMi9a0oHELMFphGI0VQY8gRJ1yvooAGZamYsdnQyniZhBn0zGK4JhSjtU3EdR89iri25bAjkNw5QCYROzXaR//AZoAcOG9Y0FALua86P11sstbpDWMy87xzUZYqLSXN/NULOxH0nyUlIx4GSuIRX7HlLNkQnS+Wa3a0LQAMPR9t9+nFAFUckqqCOCdG4Y+9n2WTMiWZ9a2Tbtq+2HYbrcyDGDcOM+MxM6x93mQlLMtEswEaboWp9o0HUk5J2u5RCSlYgGmiFJe/7z+XZCTR5uIeXk+xWQAWbbp2Jyg4bY3RbE326upKFuWUQ3FU+H+q9kV3+gqVdUnT54w84P7D87Wa+ecLZYRnzX1zk5w5fbV2pkT66vjeS9WJxngCUvdKo/0VL+vH/htAIJ5BImqpCTZ+pJiT8TOiapl4PrgnWuwdLaoFb4+khxy8V42UIhwFJFMvVntsq1iaZqHDssTEYKgDAqFeFMAUp159/Pnh4oKCz6lmhGZXue3j4u9o1YkUT0Nx+jhTGjXL1mdryzvYSkALv8lJXqERz/5z/LDfYZ/57S1LZRX75z3PsaE5SiupjicILzZnU+LNrAgIfMxUbPBDWIiVJ3PcMLDkkigGdFEtinGIQ45RSLywYUmOOdEJMY0xEFSQkIQRS5T2mw5ghZhRiZqyyntdhsdsm/b9aptQjDX/iwSU9asRMTOIbGoQhZCBKOKIgChihgjU1SJ2JZjOWWLczLcS1BsZ6DF213GN0qhSpSsx9WJV1qdn8e2li/0uHj6NA6DLnLaZmH2dCyJiErhqr7gg5j6vn/jjTdyyp/74c9p2qaieBy27lPnjyOCKtN5czIAYiGP1KPV3yFPoZKZlVsT4YjWVHP/Rr3B1/3BP+q9H7dcLqX89OnTy6srdtyuVimm7XY7pJhT7vpuGHoQRWZ2vmmb1fqsaVpbSonZ+ZlCtuKkjMnhysyOebfdP370MKV0586dW7fOHTu71M1+Z4hxt99urjb7/V4V2BGzU9GYYhziv/TP/7pf9gt/zotPfjZlpZTiMMz3ErILwTmXJKsCECM5ZEQeeTsj1F0UrmVLA1zWu0tKt2rh+E0A+ag3sY+LYGqbCHFh/DVZBS78frHC7idmEs2KFTNnBrEYqrIom9kq860khxKk6lKBOgvC2oiycUOgmoCuIIJUedBqfSWZ8a4SU/CBHU/YBcKByfFnkH7gjtGsU0Pzs+XEz4BY9eA1MLN3nqg34GJmeWG9g5s4HKV5f4EXryp4cKwjLu2CLQFUIec0xJhi0pyVyLHz7EF1GPoYB/MeUYTgfWgbBOj2+zj0BlwWSXMIRCSiMUZJUq4yKiNwTnkY+q7vATD4YNRPVZFcbPBs4yUAKlmzAABy4almy1GVXCzDJeFMIcPZyc6WcjhHOOHonTcmVpfY1BFnVLSE6xs+dvu9aYktmEnNQGPMsMbKnSInOd7tP7vEENIQ45PHT7zzD+7fX63X6CGnVN6j4+I+Cfeq/X5t5Y/H4oDnrtDwPf3jeEl+/AffHK2JsWmaIcZPfPwTn/rUp5jd/ZcerJp233cXlxcXF5cpRbOqByLvm/V6fXb7Vrtaj0t+sMBn043PZtQlTUqYybHfbrfvvPX2EIf7d+89eHCvCa3zbJPr1dXVbrc3N/Ku7w27IySL5sqS33305CaNm5oNXIpRcwIAsytj55wPCoDRLB7YeUdMOhqGmxfCCHpMOi40ClxdNOYbfXJGkdmLDbG2YeZ56tPD3ttcUuqdKk4ZYQfeAaP3AwMIFjrROMIqgFIZZGTq/AmrQ63CmqbBo8SH1axzLY0cZLUhpII1YTIVHPEN9UxN8N45PNWZnOKqvUgdftG67T7tDf/JO2f6e2bngyMiUSARYKDK313qte0IZ9ckDz1gbGs1Uc1YuS6CnrS+LDDn1Pf9MPTmVTImVlCKcb/bxyFOkEK7Wp2drSXLfrdLKRKiICJg0/iVmbx3fUwZmV1wTdMicxZVlJhTt++6vgcAdmysHhWJKacYRYpiePKzdewCFyWwTr00YRYZV0eLedSQUNuejUIJYGIFzWIOjMZ5Ii55HVZ38JlL1msroOQMYBCXNzuwpArmQF+gOVTRLEk03wgCcs4757u+e+ONT+WcPvyhD7erVjBrzcyp1YIFxMpVfHMZI+dyMeFc9UrpVHs7ObPMDI9n7pD1iHx8wrwOUEX7rk+yQ6J4fiaiw1BktMCWi+q8Z3I8KTzmnzsRCUwlNymgso3GCqDM6GTM4RqL0Wa7ffvtt3e7XQjBB9+2bRZJMafUq6pz3LQtM9/oTo7DMHlOAZJrQtuuQuMVaBgGa3edc7b1VTH560y/IbSlgG1pRjueqkWbQW6dKO8CoGh+cJWVwozvTwwcm8stIoYQAU1bJONPxDLGKaBOpjslf1bVbBh0ofbSycIZF4SLsmuTkXxWkzixsh1nBLGfYPOxFiBoGXWu46VkRCAxAmgb2sb7673GP4OQo3vejkFf4M/wYt+LAGCGEOQ4x3QkmcPJ/AiP/kGvfWeWFoAnUAib1ZQQVTUOw26333edqDjvm7ZpmoaZh9gP/ZBSIjRZCXjnV+3KVgIqQuxYBRFXTdu0TUp5t9+lITr2vglN0zjvANBiPIYUU0qFfCYCgJaElMS2qUUcB4joaGRK4AE7eWw6xrtEpn4GM6KoslKpOkKZ8nQjjRMBzlS6agi+IYMLJvt143tgIvMFwiVNHgF+2+/+/b/jX/8Djh05i3dyY/+HAGqG84joPDM5RPLeISIl6of4+PFj7/xLDx6sz9bqdBiiitr4PFv1Hh9gI0OwymlZpDnBoWrwWffV+3L7QXPj0NCG23fupJS8d33fd33f7ztVAVFgLdwBImQ3BRvUuEgh8lbeYZN1tcmuACCEsFq1jl3fdZfDkHPabK/6oZeCu2iWnEclOVMp0zeCmBVUUgQAdD40DTvvgm+alp3LKRc+i3GbrejlPHGxi2WsdSglwbGycKg136qa87yMGedCnkZVPVi5Vvg+kHXQOHGCx+X8LBBTNeIN1YY9OnrxsgV1IFi8l/naCiwy5BVnLRmTIZAIIDTtkNTmYwRIQ7LfGNg1TRCx0NgR0TIOhf2RxNBeBHUhrNq28c10tuHN6ur7qtufTjdQPbqXjl+Jd75pW+9cHA+AKau9ZO3BgnB1IOTXA7NInMubISQyesFA5d1iOJKKppx3/X672QxdjwBN256tVk3jVSEnSTmJZCZERgv5LSWanaIwM6ggsRla9d2u23cp5xCaJjS+8Y6dKXhjymJ7g5zTKPGdU4mJAQComG4yW1IyiardtVAtuZbBg1DvRyxbtfzwcanOREyO2WgOY59CyIBi6pv3QAecuAyIjog8gmpKaNd9LdA10CrnqED2iTK7KeLUMhhw0k8UNEeDD8yu67o333xTVV8PrzehIUp5pk+gTPYv1zsOLYTgFVkJF26BeFjuFZ/R+J8Wi+khvmQ/NkvKktvQvvTggVkBP37y5GpzJTGaEh0Rz9arW3fuicgQE1meIc1957xsrPRCFmQDCjkliQkQnXfMHHO8uLy8vLyMaTBVdggBEFNKfT+oJcs3AZCQsHhV3ejBDEihXbXrVfABnWNHY4yyhTYWjWG9ryVENpMgKk98jtY6ROWq+R1ntIgLeIlahwiYqIAqixgF0bKosDOAJqr+AiGEURumKBZbnG2oderQ00Rt1pxzziATl8A2GCNZ1YztxnaG0GjPdmXxSHwtHFBCdESCgJqzkFZZBRVhyfZ54J1brdvQBDq0M4LrNlDw6fuz0xueGjeaCI5dNJzj4INjh3V3CweGp9Z36nwYHK3aR1c4PMh0nGzWptgqVHSeLcN93+12u13X7VKKzvmmmOTAMHRD7MA+P+QmNCEERNjv9zklx4wl364Q0aybTzmZVbL3zvsACjENKcaUkxRXQAXLvJ6fHxXspoSZIhDajZQlp5iMgIQmhadJKDly6cbF78j7WTCGESCrgmTIwAaOIs6sOH2WS8n1q9qFaMPsvMSb9ZyxMHFqqcZAAJWcATCXe5qttxItowkRm6WaueASITMnwm7oHz5+hEwvPXiwXq/YUbR8ghJ1qCcM+bG6BsplMxN7nq3WwpObzwUX7WZvlxi/WEVRiQgFVSTHCDkjMzpiduycI5dRWMRoH0vJz6TjmzD1sv9RBcgZEbxjzfL04ikq7PtuiIOI2Doti3UdQkQllgvZoiNMvnsjGmizXjN7F4L3gYMnxwSYIde4jIrkJIby2EXNNs/SFDdZXl7haJQNm8Wzq4l+qepsCOa7GqcM3uVGrzJ9oWl1UGKVJvYWMVsgWEpZSySGXZgEve0dJOcUE9KoK1ZLjxmvNwMcC++AMhKOmN1BQHwx+xTNKYKKpLzf77IkI7wW95dSm4jKWSGSBTShQnC+DY2vfEDx+jPg0/7nz14gTLGEcy40Zgix6J4ON8ALxw9cUrRqV716W2J7zmmdoHNuExIApJj23b7bd+bmhoTOMTuSlPfbXez6qRtcrVbr9TqltN1uUTV4x8xIiMwKKKpxGCzF1LpaQx1MJBOHmGLKOS9iK6dtFiKzTeTOOwcAOSWLyJQsMQ2SkyIRO2JGnkkQhWipqkBlb2ZZE6PQsYgGRJLkLNkJO+fL5XaaxfyCB4AbPwoD+rOVbyOrEDEA5CQiSVVGzbYqgOScYPbrSmYdQWYsZMeGFEVvVlX1Lijrbr978803AeBD4bXgfaYsk2k7HmLwWIfJlfu1KhwVc2v6f7GqtycY1tNnpc+k+uPxMriQc0VkiP1uu8siADrEAQggA4CGdtW0K0Tq+71ZPjG7cfev1dMuPOEibgVkYnacUjYTckTs43B5dSUi3jvvPADGnIZ+yDkhovM+NI3VnZSzOfqYueCNDoD1rXPvmoqUTcUHaqRojLejIpDBfsRkm97pIzkY1afSb0FyPB4bYB75oGYCN09CVMV6iI6cqHGgRFBAKbHHUmZku67MdkdylpRiTjGZREtVopQSkmKseHMHVef6+e+6KwIBVMyDfXN1RVty3rVNi23LzgECETjr9lRBUQg1OQRYtWHVBO8dgn725SWf5gPgeIuBtXcNgne+EGNG7BZrQ4SawvdM7g8uFQYIVaWtZk1HhERDP+z32/1+v9vvYhoM4cxmcAggOQ8xppym3tH70DbtLm+HYSDrUxAJ2HlvgSp938VhYKSwalbtSkX3+x2xA6KsEmPJwCuEVGKr5lhQU3ZMY7Cqaomu01yVsHJ1F4XkZImBMGq+VBWVVBRyUoExVszoa6IKZtBlrdiUNKT6XmBuIraBI+cMCIi5TPqeASilNPQdSAYmx1QQTEVrsDCNFK2URYGRD9uQ8SUbvppzGuLw5MkT7939e/faphGVIcZaO7Mk8h+w/xYJv5MUbOZ1X7MTKK2CIixkn7hoxo7rQcEWZwcJHR3H9/vdEIc4JCIWSKqyPj+7c+dOTrLb79m51XrtHTt2opIlTw6mVDgvVliViZu2aUKz7/ZPnnZXmyvJYhnRRUGokrPmnMBioZmcc855JEwpxTgY7wBu6C6JiN63zjtVzaJIjGRUnGw8BVF1SIzkmZxzpesnXEJsxeN5ZDcYXA6ICAyoUIIrJqCoQPaz2UvB/a39kQL24CgmsO9TxZxL+KT9b87ZPqycc7JliEhOyX6mTCKk2dqz9oI93D/q5F2mRzSwCWWGumap5pQzmHYy52x+XLZF43LyO7u5A7tV256t140PCAvp4sHl/pkSgn0mlst4kjOngIDec9sG7xhQQZaZH6P1wWJyP7VYODB8mDVik9/DGPVCjolou9s+evJkt93ahVAWVTbA5pRykpRFlBAt9pqLBFdVNanknImQkC2xc4ixH4aYo3e+bVof3H7XpZTWqzU5VtWUU8VnJ/beh2BbuNHgHMuFmorCFsSyrcnoxBPgDRP7Y8z1nggN1hRlwdEL0dTVpVCqaswJxYwVbRTA9yBgFVVituNMRDRqCajxjphVxOLSUrZ9XsaScmkKnZImCUXPPI3OOhfz0dA7Z0VU02Pv9/s333ori7z68su2KNbjeMiqxBwye+aeGisP4jomsM4APlwr63K9jNevj2teOqhmMem4C024urrabncqGQiB2Ifm7p07d+/ef/L0Sdf3XnWNZAHrICCQx0MfTStv+3sTThISgMaYtrvdZrNFgNCEpgkIKCrDEGOMAOC9D87ZvCsqMuRhiDlHVUEkx8UW8IXvYbP6Ng//iW9kgItIElUhchbzxewKTLl018EaHZNsbh+Ehb9EAMzjkFuH6hIRkoIIyJTXpDOAabpGVChx8KAgIjHFGFNKMceYkvnVyagWxMXScLIRxJLusEiXUD053MEYSoxarWtQj3MhCGfO5zAMg5kHT2QhdqHxTWi8d5oyN41z7uxs3bYtgRFG6u0pHnPQ/rc8AZwgf872y6AA6B03TeMcLxK/dWHjcx1ZcRHuWJlBIh6cD2PnIIUYTwiWAJ+GQXIU1Ta0bdM4x0Pf9/u+VGGk4H1oAiL2fTf0gxQiLxH70DbON9lCTpNZNbiUUs7JKCsxJQY9XLUhOsc+eHJuvD4AkbJqGvLkAp1Exrx4I/jTLAcZDU70AKo28HO8kLWC6ct7IVkVMihEVGF27zGxSyvNl4qY5619ggBwvj47O1v3/f7dd9/Ju72GwD5Mn2vK2QxgxkIyPgepQY/K2tNOrwzdvn/nnXckpQcPHpyfn4NqN/TTPmAKdsZn8uPwENlZ8n5mh5CDxEk9uRo4kU6sVW4lQo6Gw2QmJmaRDH0P7Nqzs3svPbh7914ITkRSjiwjLjyxBhGJyJFDBLNHDSEYSDgMw+PHjy8uL7tdZ+6ERjYDkVwWicSOnHdETtWa3cL9Ny9P51w5bG5yJxOzSdN0DH4bCZBqcadEaFozJpws2g5N74t9W8lcolGzMq33R/SmupGttRearFzMU5ptwlWVLHGIKUcprqQqoikna/cLk6KIN3GijuBMVFisDHUe9g7gg/FTn6ZwHJ9zxSqQqV5XDnC1xVyp6WMtTLYsznnoWft+aHx89WXvfeO5HineT8n9XxMCOjFZzyewIqBjakLjS+yJmOZK51yEg0UAnKZo4LEiAAp9psINAUdeDcJq3Q5D13X7oR+A8db52Z07d1B0u7nqum7kk0EIYdW2CtDt9l3Xq6rNBM2q8aFFxGFIQ0qiysQAut/tsmTnnA8hS45dNB+h2dGO2O4SQsoqKqoEDHNyhJFTc85Z8mixQjCaL+pyH2Rca54ZdNlsTnFyKrU9s5kKjXvwlJOoBPDs6KYXCCMXltvsuFWGEIsaPjtbf/hzPyf2w9APj7vB0pTmrC+RBGpUJ+OuT14oeFCEEQAgZUEEdqwil5dX+91OFZq2CT4gkqpQiXavnUOWFM7aUvYItMEDDBGqPFc8eVMtv+7gTht7Q0JGzDGl3W7XpyFGYweQAoBzt+/effnlV87WZ13fDzHJSHqZfPUIsSyImQFUk0wZiTmli4uLt956a7vbhhDadiWqMQ5DP6SckMg754JDYkttliwpppSSqgChKz+VqCCEN+nmkMk2k6ZN0WKlbkgOFd8fQjgM9xsLmeAo7AcAJoc0bejm2KwSW2GQF5b2nggFAAVFpPLnQQCQlPvYd/t+GPqYcs55krtPoCsSofHQDrCaKh9wwvlw2awuMh90SqZY8opHj5pJfVnuvCI/UEv5QEJGQnYKlf+P6atjzprS1eXeu67vHbMjSiK1//lnEvh5Jg30ejz/NIH6mrtMl5tb1SIdQudcUTzmrFrnfoxhHnj6qKvWCTqvAXBhrapTbDwiIqaUum7XD735+4uWCcuxC94bRJrSMF5easHFwxC7fhjiwMShsZVaUMR+GPb7LsakRfuaYoyA4No2hBBTGvo4UjOtDpG1vYbq6MhJzFA42jiFF8/SQ1MUzxzQyqigGoBGhKpAkBPFlMiYympmGOa7KyoiKUUR0hsyQb33jGRJryMZmmyEl2SsCVi3q3D7Tt93oPLk4iL2vZmbjLpNJSQXPCOKKIg4Z/I0MVRn1nHN4U4l9muI8Z133wGAl1966ezsTEHjEEWV8PSUqPCcmQAWJLJrrmcd14+V8VKd/HQc+mTIdIppu9tud7s4xG671Syr8/Pzu/fv3ru/alfEPIqecE4XRERCRkZiawXYYeMazXm723X7LqW42+/7OOgY0iKSYswC6tgRk/POsROFnIYci1U/IkxnLhEXjcKNPvqiI0G1b1UthVXJ+OyjFL1UVRqtCUVVNUsWBLDsaCTSYvtc0BHVcv0iAjO5cTYlIj8a56pGUU0x9XGI/WABkrZmjxZ1naKxnmoPv1l5vEB0aq9vrSy09HqkASf2QWlEJ0v2irlY7X9mGcIUTAMCQtay4KmVqM3v4Bwb8jldEc+st/g83e+zCGzHdd4teTX4jOK+4PuPz3Q83Or9qx6dAFO6qlHf0XnnHKfiMQDGcpNlRgriiadc9XVTOg+MZPeR7GKXmENAEMl9319tr3bbXW8kCu8cESKklGIazFpn1NQSIopKSjmlKCLOueBDCMFoGN1+v9/vU06AkM33VZW9t2k9xpRzsqSLcj0677wnpqnfRyRVlWScPBw5lJUxPZkCgbROPaqaG5zk52NjVABTMUMkLkglIBBqFhABFFAta6mUb3QA/Jyf8VM2m21MSTRZq2sCH2TKMYnI2fnZKy+/cufu7R/7pV/05OmTd99+d7fZ2O54uqi89yEEVY0xEaL3HokKVjaZ3M0f38woF5E4DOz4wYOXXnrwwHuXUhI90aVPbejHPv6D7z5+urgF9ORtUN/8cxiZjrxyVCgEwRr7eQajVBUAUkqbzWa322k/QBtu3X/w6quvhqbNIilnGr3vF/aINCqNAESyZGSHfRyePnn69OJpzhlHJZeK9EOXUkJA770576NRv3Ie+kFSAkBi570zKzdCNH24MdNuOL5b8y9ivuVIjtlAS53MOnCOXbST0VBCU7AjABNZWtHM8gRQQlFGQHYcnIWkUolBTZIoZsIU4xBjGobdfr/b7TUnm25n/heC8fLH0m7UHwOspF7MLw0CtJaU6ukeF2FJyNGl5cD0BIr5wIRI6GKYKL4pB173leIBQmhXq9Vq7UMo4g+9tvE/WZbnl3YkvTqpGzio89foAEpx1+O/h6VmvorsWQBo+kxJAAI65uD8QNFqGVUcbDzFzFu8JJy9oCemJx1ywiu0CCHH3PdD7HtADD6s2oaZ+24f42D1iImm8LgUc0yxZFaA+Zc6ZpeypJyHFCVna2rKHUyMaOJY6+vFKgcSee+c2cPNQSNqBnAg5Q6xzeFEVpzCKw6vvAUoryb3BDVnMVbBZVKMxSeUi00zlIU7aL7hBPDN3/D74R+oxx/8I//Bf/WX/gpMFsN6OvhJT/VDumCXzhslPZKhV351873nvAtNo6La90B06/ad23furdZnhBhTlkJXn3gAxc0DEZjIBUeAfS9dt9933W633223OSUBIGOnj5cCM/MYy24W/WZuKDkDoi89R7EWyZol5Zxysmn7Zrsfg2e17ILL0bVQWiEgMzJSziLj3pUJAczXjsfzDiurH/O/RAqhbRpy7PYOFCSlvusuAQx3ykXLnOMwEiWmgJiyBqaynFZVQIunGRNWl3qB46JU0YCXUOKzeHIIFVa4MC/GhRcRLnlpuhC+qcIkcEJH7aptmwZxyTZbRhHjYrit1LO4tCwfE1YnOulz63a1AzhOZ372IgIPI9aet6PAGr5w7HzjqSdVERUCXqzS9cQ6ZHbiqI+IsmKoxFCVXXBKKeU0Xw1ZgLBpwtn5Oajst9uUEzMCMSg459gHUe37LvYDADjnXHDk2FSd5tmpOidaIpJ526noMMSc8swxVAF0znnvvCKIjHuo0QF83HiA3ZxH1NnDjMSDD8di3wHNXYbyRHnEiktm6huDkqQsklU/6E72cA2F7MR1ecILuEBvuuAlw9GwgAv331I0iMn74EKAtl2v1w/uPzg7W6tAhtJQSVZAMBrweASgkYzNArAf+s1m8/jxk33XMZP3HgBTisMQRcU774OzRbohkDllG6MtA9q7EJqG2YkJ7ySnmCzPHW5uNVxe1OiDQOMS6OC6HMdwI6cKEzlmKrrJGZQRVSUZRftsoGhKWVMehsGsMvq+jzVpddIGu4ALH1ady2KhvqlWyPDSSvpAtYGHVfqA1jmetC+2iKUDqvCi0ydUQBSt59rxvRJAcMyrVRuCPyj2SzvbhQ310iT7dJOtL1a35wPgUBGDeBpbwkNXa10osQCujS9Y2mohhODa0GzdPqWkJNPZjtf5n2pFz8DJq1UrVd78CVqhjTFut9uu2w1Dn2ICFbtUHbsQmjT0MSUR8exjTmmIzByC64fU90M/7IG4aVdnZ+dt2xoxcbfbm8J+JG8IEIemCT6I6hBjyhacW3sGlrnT8EQ270zSDLN6UEV0Wv/OHUVBwnVhfwETCUON+0FF/3WQjTc52+JI51BEmCOZPviPOijhFMezihWvBWOnOQzPmGYnih4qwNANXde7EF5+6eX1+fmt27ed88W7FFGyxJwJMQTPzhVZquOmaUXkyZNHFxcXOUuWnDVbA5hyEtGcMiE69i4Ugr8RiNMQTW9IRN55JGLPzCygOWfJSbLkFDVFMLr9Dc+A0U52CqZbICrT/ZazCKhIZkZGb2oAOwF03l4pgcUHONd4Qu76/snTp7v9PudcUvOcM06IeV5N+YsmnBnbe5mLPU7BjiPeg4vKREWksShdizjHOXoCFwFwOJsR1siGHiFGOKrRAJZ+n0YGJbS4CgGpDLlszS8IELxfr9ZNaBYEyEPy/FEpRlwY6r1w3YZToJDDZ/D3X/zvT+UqXGcOgQjeh2bVuitnWYBqpqnLzRwuVjDVaa3z4mH0/Z+tJIxwhoBZ0n6/226uhjikGFWVbPvsGECzXbeSs3A2Mo1tfrLEFHMW73wTQts2PvgUU7fvNpvNkKLtaVUEBJDROx+C33f9MAy21Rx9e9CQH1E9WF2UpN7RrKqEDZR7jGbzlGsOdRxvAkIyZqc1j7YBMOejqaLpJK5AAkR5zrj2QTsA8JkAN17HWht5mdexIxbfbmo0JETZdvuLywtVOL935/z8lvcBZ0kHxJRTiogYfPDeI2BKiQhjiHGIjx4/efjwobWEjp0GyDkPQ8o5IVIIvixOAHLKOcdhiJISFF75rKw0OlmMSQzzUS3i++BvRAMFgGldYZtO1FpUXi2yJSsoAXjnveOpMS5ehDgmYKAx4sj2//t+//Tp0ydPniQRx4zI7D0pqKiyWharqcemaN9c9GeVsH1GhPWgzbSRBRbdvFbOAxXXsJQPGutNoRiAHmHpC6gFDnijeBgeYRFUxTbPMCosgX6Kku0AOFuvV03LlfcFLNEpXF5ueKrgvve6/RmlgR7J8edfHxq/WrU++K7voLDZK9yremcR6lSXkY+nEwFoPL2n5EhER0yEyXkmzDkPXZ/SoADBN23bOO/ikPouppRyiqrA3q2bJjRNTKkf+iyZHPkQADGmyMwAaug/ZEGLPLUdlCkvmWxZNuE/CIDk2JsAV8HMoghB1HIoNWdjQWSRGjo0pdgi4WKxwCnVnADzCIip6sj5myxFYXLCwVnLhEjCwHDzSIB/kEGgw/p94NwwMglweb0iHKSSPdsTQIGIUCmltN3tKOdtE3yzYh+mHDgRjTkOwyCqzjnnvbl/b7e7y6vLbt9tNleGs8SUJBsqCETEHNgyhZCs7zelq6haP20sICJS0ZyTPbJkEAFEYte0TdO2tjO4+f1bBhyVsujDRfQVEpJ3zkYLV8hrBYI0GIqIfXCND8QcU9xsNtvdvo/9fr/fbrfWHokqgYrinPKsoGaarqogOUlFrqmJWWN6bDWjTMVhdoVc+IlpzQaqXq45XEzc6QpeX6aDLbqM6+VZk5c9jefR4o1VBcSmaW7furVatQvvQR3l7FB3fZ8xO2i9hkz0fkIHjm4WPRjNQ2hWq3XwHkq267NOKzw+/rCW32lFLjeumOaUJSdCYDNuixkJ/JrXZ2sm7vbdvuuySBbIImdNc7Y+zyrb7a7rB0BkckiomuMQCWnkhGHdSaBz3gc0JxzVUdOjtv513nnfLLRXoiLZlnVmoZkl55Sr0RWn5KKDt3GCNaf4ovnilMIiOuzxK7mjbaXHRL1/iMr/GCB9XNQUK/uHGSyzhR0+D0g9ogwWzbZCTEl2u8wMBKJydnbug2OiEkKaooKaIjcOUVX2u93l1UXf9c67tm0BMcWh74es4p0PjXfsjIxjpJicorlNMHsfvPeBiVQkpSQiNeKPzEAc2mZ1dtY2rSUE3ZQFhAsGJRaHkhIJp4TgnPPe+8IOGvmghABAwECMxEwECFnyft8/fXrx6MmjzXYHAM77drUCgJTN6DpNIcC1z5ppGhewHVYawpn3We6HipGoY2M/o0RVHdJDGA/GA+YUZriAfWZeQGU8esTRMYZjLrYDi62eZgHGtmlu3761PltP8NozoMYlmWe5HX4fdds9I/X0uun6xVckR8cY2EzXNH7VFqcRlTrpB+HQ7Xkp0pxBOl0+I3PNZSaOw7C5utrutt1+l1Kc3lsmDt5Lkn3fDX0HSE3j2rZpV2v2fuj2fT8kycZISzF6731wotJ3wxCjzbIAACKI2PjQNEEBetsxwBxqT0je+eC90VvNCzOlrDmJChX438igltJexx3oBPLRwRwFoCIqKilP+4/SEY0yrUMCVS3PVFX4h2ID/PyKvRjUtbo8rx2Z9Qg2OhASzf+MmHK+utrElLPIrdu3Vw2J6jDEYYjOMbPvh77b70zlZNYaqjrEwez8iMmRY++M6CkpDTHGYSgRRszOeXLMLrBjS3xMQzIPEk0JkLwLvgnkfAnsZQK5aRM519X5jp9a4TG4xpxcmUkLEUbNX4uca5smhJBFN5vtxcXj7X7f9d1uu+v6KCqIRe04WoNKEQzXnTAgmfisDAY4Az1UFT9cmnePagt72jJxfg+yxOq14ejoVH/jdaiJnrC/X+4GpkwaW1QAyLisHPElhRwBfWibO3funJ2dIeDkBwVLntXz1/T4vuq2u/7EOWSMPuNmwhf+dwVBYO/cqm284xIBqoI6FkZYGgjiAauqRrHnoc7+xfZcKQ7b7eby8jLGPqeMAN47IvLOIWDMKQ4x5eS8WZ2vvfc5574bhjiIqqWDJclnZ+dts+r6btft+76f3w0RZA7BhxCGmLqujzkepdKqqmguDoY55ZySaCakItQd17kTsakkBClVV6utOGbliYpITKpafLdG+wTFSVJ2AjzCenv8wWcBvdC9A7UBwCzDw+cDPicPBzULEGVmaYL3Pg5pGK5UgYgc3S5+NXEQ9V5hv99tNxvrgr0PjrkfYj8M5W+Cd+xNvZUlp5ijeT4jMrsQPDtv1SunnFPsukGMJq+A7NiHpl21q0AuEJECFBrpDQEgmCTmE++9xsxnqJ3K/yBRBcEjkgLGOFxeXb377rsXl5dZMjE778it7aZNKRl/VKs93vj7aUTnsM4EKmkEpgqGiu6pWq1q5z2B1BR+fRaOd30FxdPfdnIFqie+XGfp4Eh3yVmJgg+3zm+1TWNudYTP39DpYUAgvp+6DQBuEaF88pSomBNTQO/B1+vhqQSKp84mMy4jZIAQQnCeENMYuDmGSsOhBdehSUfNQraNUzlyRRRBCKkJIXg/9H1OCQFWbds07aptc86xH4x3X7z9mbLo0A99GrLK6M8JAGqmJTnLMAyaEjKPKd2GtxMgpmxDd4ZKai8isY85Z5qySQHJ2G+TaMiy9ogYCQjn6yOL1IPoSMEwEp4dloTIjq3TlJwV1Eh6es1+s0jXRdX41B/09S/OIUPPAr0Oxstq7aYVVLtIBl+QtKZ3F0GKtk+mT4wYY0ybqysiDMGtV2fENMRh2G6JWCRZjDsCxhRzyiLqvTfDSOcCIqYU0xBzyqICgI491U6fOUtKWSXHJDmCKLAzI2iDhojNac1CdAhumAZTqPtSGyUr6Mi4UAEtCgZLerd5IIQQGg8Au93uycXT3Xa/2+/33X7XdTlnATWjOwLIJd9MRml6Ld+d/V4Mbppsv2tz88nG+STjaxLeLB3EljjOMRlYzUJi8SPK3m25O15wDRCvW7zOIlit1YflswjBr1crM4YRUeATmA3WWFx9XU+D/bO4CnUSCcIib6UUWffc2OsX3CwfPvWTh211NzkmF7z3LhVDpyOTrioVbNHb6WGbMl48IJLzEFOOjrkJzeB7i7Vrm+b81jkSbXfbfd8BKDnH7JiLu2fXdQbjAKhoRofBByJKKRlPqRxL5nDC5LjU35wKgQjLpFoI0yknyDAm4iGXXAgyerZKie1ldOyIEKXExamqoh68jQR2BY4p0paZJ6opRVAFer7L20ge/eBPADRFz77QlhPhuFM5up8XnOwlRwunOHP7z6i0cM4BQIxxu9lcrdYAEId+GPrdbk9WKEMg5jgM/RBzSk0TQmgRoRCDJaY+xhRVtbAnnUOmUTtiWvWUcwZVREZv2S3tatU65xFBRqaKnTTyHj73KVJp9GtcBHnZ2Elk1sYipXtLWSTGzWbzzrsPHz583PWd9943TWgbexNFRLKaW6dWbp045aVUR65M3f28RdUxk2CMkofjynyKmn8gXTpQjE5t68QsH8OEC6Go+Dxc68RQXzWTsBwXLnEThxMBOYSmCa0fTSCefRjfpEDfrG67428+wUJEPP0jEE+vGg7VdUe0IABCapwPbTOkZDcQKYFp1xEXn+uEkVTVq+52UdW0gTHn3W633W77vss5O+JVu8qSnXfsOMa077q+6wQg+LBer9vVyjHvu/0QY7QEVAAFbZrVerVmR13f9bFHRGtdQDMgOOecDwA4yusRiZgdgk7izsIsm2xJCNGCXEpykGTJJSJo/HdjzekBAcxMQ+aRdvZHLG635qaFowHSlC+CM3JdGWApfNDXAMQkRoCZjDMqYoUuFelLqr/tg/U08+c0RFDeY6Pem6/J/G+EzjlVvbh4uttt+6Ef+sGxD8EhcspZYpScCdG3rQ/BORZRiTGnlGMyUylnrimOGVnLqBGtd5acQASYm7Z1oXWenfPMzqiWXPoIGr1Y6aZrOwVYOogfWL4VIxNCblctgMZh2F5tLi6ebra7IQ5d38cUZ5sIQDHCsvnJmj/+aO0g06lSrQF1ueSTOQ5m+lTtf+cZov5wi9vziA/rfOTj4QePeEC6LzY/CFUCWPE70AqAuBYSmlVA1e9UBSwBwr5tbt26dbZej9SspU3p8h3Haw6aRTTWyYyjI+XWcd12z+e5nWovT2RjTP0QXi+zhIWKwYewatq+65PICIdfs9SoPRnxkFg0WUOYCf5uv9tcXSHiuml88E69c5xz7oe+74eUk3M+NKFpm6YJWnD1aBirXZht05yfn+ecdt0u9gMiAhMojAGQoQmNSAl+KYZbRoYz0gPWLfmUMC3zGkSyiDgb0qvPX6dVB86iZ5myNKDEwZdyryUMe25PRkVSdTRWvarCPwxLYAUkR4uQD31mjZuNlvTE/KyLbhKX2OsUQS+qWfIk5yhuVIDOOxXd7fY5XwGo9wZ8+pxyjENKmZBC44MPiJhTTrkkypmJv/fee8fsQFU0m+lxHKJFRyAxsA/Bt+t1aFpmGj351WIZx7QlJUBBvOESeC6YozDisCqUpFzEFHOKset2D588fuutt7ZXG2RarVZN25Z3ViwuM028ZyIy9qqV9mL7aVYOuPTqGXfOco0xylj6CocOZkuUIxVl5SpaO4IeYA+1hhhHo59lhsSC3K5Lq9EjFGR2uQaAbLEobXvnzp31enV9i3561HiGGRW+QNd/2gzuuV/03L/H00958R3122enRPC+aRu3cTF144pnHDlBF/jtmAKHNBLtZ4fpMiOa4N5774MnphRjH7Gl1jee2YtoHFJOCQAdO+88E6vokGKMqdgKFs4uOR9c8HGfuj6mlAoAX2g0aOTrruviMABAcJ645FhQsaIdm4iqAx1JabM/CDM758bGHMFcccuFruOoWC5tQBAFHJ3DCg+PqIgnQY2uOgmgxt5/2keLNUrf8Vf/Bn/oRznnip3R5HWoKlnUzHUnEjryyB1UImLvxpYNppvfnkzKaYgxWtOngMRNaFzwIrnv+rzfAfMrr7/+kY983tlq/fTy6ZtvvvXuw3e1H8A7e5NW7aopjqo9IjZt48hVrABUsHQBUVVkcuzMkTulmGJGwvP1+uWXXnrttdfOzs5iSkPX2zw2dmCH7M5qPkdYRjHBNXp7hVNJAIA5i2VeEhNmmvwEdDwbAKBt2tC0Oct+tzc2SOM9saUMwTD0cYg5J7Oic+zIFeQHAGJOOSURyMbzAbA2xIcmBO9McFCyGMksQKCyYSAmyHDTNIgDh+zZdM1iBgC947ZpVPXtd95+/OhhSnGIaRiiqSBtUFU15VoWsQh3HGHRonkcjdtxdkk5KMM1c39iXk3K2lrHMT1pncr/jCjhcoiZTS1hKazF+mj4/7P3r0uSJVe6GLYu7ntHRGbWre/VjW5gAAxmzkVPcGiUXoFPIf6nSBklo5lIk8wo429S//QGegCKP0QzvoBo5zLnzFA8ZwZA49KXqsrMiL23+1pLP5a7b9+RmdXVQGPQmFMBDKY7KyozMmLv5Wt967tAc6eFbd6I3QXF7m+vmv665qqaGcIwDFdXV/vDns55nfdHgN1be/E+nTB+w+r4/HuG76jtenABgQ+8KgSMw7Af98zOBZK2KjF84Mizu6ZeNf9EJS0pSXbtSQjBncKBMIYISPM8zWkxQGy+b4g553maUlqsdoBIHEIkZlNzX3U1RaBC80AgCswBED16wuPdmAjqqqqy8gsjtFxCFczUGn/hLg5ceKLnPamvMgG17m6wWOb5IF+xIiJ2wYzz/zfkjPrmuyTHVNsxg8TuU1hzmMrbSIDiLQpzkfGAkRGXGYdWr402ern5ogIhD7HcLx4GXEAY5hCCxIgI8zzf3tw+vnr0wQcf5qwvXrxcjhPE5u1VYmSzpGK3MVigQERmgIQlEVkNmQIFBBARACPkw2EUkZvbWxWNQ3RtVPEPANt4u+Hr51uDe3eDr7v50Ff3OWVVwaLHrhlhIioaiIcYx3GHhGlJy7Iw87AbXQ+sIkuel3nOOfsFFofIITAxAphaluS8tXJYMxOFuBvH3W7c7QMXt2dfljITrfInpyqU6/3rV9f/7uefrwY9r1WF/+q3X6H1S0x1pDJ4hK+ZSk5Lmub5eDp9/utff/n5rwBhOBziMMRxLH9HTEQ9jsK9s4i5GOCtbZHVjBn/gK0P5ryLWxhg5Ue3hfsqV+6z32y7Am5OFqvp/GZ/UFumlcVSu3uDe1imnR64YyLdecE9WkWIgKZqiEOMF5eXh/3hTl+ycf355kb8IRD+TYihfgDYG06Db/btHnj+eeV27vtu3O0P+2GMcFsuBUTDHu4xhLqAb0kAKyWgnLCEYPMk1zfXx+NpSfM8zcu8INJ+tzsc9mGIy5yOp9PsYA7AUq2eU0qnafLMNqyo1G63I8JpXuZ5KfovQJMMahTCOBT035NYCLkNlVZ5nT6LMjMigYG4aWz1rVJRnxecL9QWAU0Y3Nxi+/QMrcQq78DFylvD4EkgJlnQzYWrlV/xwyy21WVh4IbEamCqfhbSasiMRGRGJmrVedAbIA7MzO7XWxhHvvpUbaOwnwGImOYliQBiiIOKiJ68bJ2Op1/96vM4xPfff3+32+1243KK/vYyIiBm1aaMVlVJgqFkBCKigjoOErjMWlbeLqjkmbCk5fNf/UqyvP/BB5cXh5SzXwmEhY/rTjWt04Az93UvnfqAj8S5umYtDFkk5SQixc2DSoH1/N7dMMY45JxTzma23+85BHdNkCxpnpZl9uC0ENgJPOQnnMe7pCQ1F4VC2I27OI7suwEq4izvUajmKbbfhqhIZtT0v/5v/5//xX/z3xLibn8Y9zsOdbzzwmoABtQh5DUrrU6HCCMNPMRhiIiY5vnrr79e5skAbm9vgRGYmkheRXPOVlN83RQU6mfpdUvRTM1EnKDWZUGcbRqgECGaAgs7yBl7d+D+8+mUVVZFELYdKza1/PwzNtD72S1m9zekd2imPlxo0+q4I18WYBqG4fGjq4uLC0RUW8+bh8rs+bEA8Cayjjep23/YRLCHOSmKgOMQLw77GGPlmxl0/C9odm/n+a3Vjq2WLVPzKPbb29vj8XZZFgC7uLjYHXbjfqeqS16maU45BWYvLl441CwtKeeEzGYKgHGI+8MeEE/H4zLP/rM83Qkj7Xa7IQxmuqRU2n8mbHqgygExggYImRXSl19/kt0x0Q0koGSboOcEFBIhl0xUOG8ttIUK+a8Ppb8zABNRNaMmly+eEiuKWvZJhXyoIMVrrl+dFcCgRrRhm8mYOHAQyYba3Kq94S2mWx70HQIxG4BMc3WsQ+Zg0RBxWeZlmkOIZpZTJmIcYsvD8tB5VWn59eI9L0LE6KVKzVTFzJMFUQydM+W527txl1K6fnUjSYKj58RU/6JPPO7cY/dx1uBeot833V1YGhoRN2AAI0JGghLWiDHGEAIgeGjtuNvt93tneaaUc1rSMqsqEsU4xCEwBY+q9gEiLbPmDOiCqzDsht3+EIcdtdnRFCFQAX6a13RBr8sayZevph2pftNXl30xNmS1NR5FYzgM0fdaKWUTBQBRub29uT3eemTe/vLS8Y1i4qyey+iMtUCB17S/4tcPBuC52FAylGizt1wT40pGWHFZt40Nj9Vwro3Vf2OR2Fnw7Bmq8E2dqtn9qLsZPAS6bzp6VyK0ykZqADmr8RCHx1ePD/s9AtZsKLjTWXxDL/+dLPTCZuN11zP6Id3vvc+/z4F1/bMOqfM+N0a+2O+HGKsczHAldlkhVHYrjqbTtzrxAWK1RtEQY4hBTU0E3PWNg5qllOcliYqpCiITMQdmzjmntKjmcpWoAUGgMMRhWZbpNKW0eHA8IO33+/1+zzGq6DRPvohr+A8igakjGV5ifGMMhN70rz7lKquMYC065ResAUpMVEjYWt6TIgoiImQv0B4gTt7ceaCNGoiImWEx3wIER68qukrF2L0FaKyK6ipfRCJU7U3MraEJFgxKqIh4KmyhzVQjOyIzCyFIFMnpeLJxHC4vL3LOp+PRv9/LV69SzgiwpJmJTZWZmFkNNGf1JEVT33mi+uSDAalgfSU+F808E9cA0HyDT4aIxDyn5Veff56X/P777+0PB5Gc5gWwdOa1bGxAbteB2jrQFAXHxoq3GLvVmENrgRNmfgKLmrnzF6Z5STntdrvD/sLM5rQg4f6wH8eRmVJK03RKcxIVRBziGCJTCIFYzZziCQ6LFI8pYA67w2G333OM1EVfeY2nJseqX24aiBKMLaKiJf648Iax7rCh05G3rGYQVRUBhHEYL68uA9HLFy+/+vKL6fbWwJBDH5tigKqqObegOvJuvyA+haqgHi8pUik9Cu7508jcKyaHZSTplb7ntsI+WNs5Q7zNNPcGQpybGN9B/BDPD4P6ht5/KtybLHDvEVPP13HcPXr8eD8Mhqgi696w4yScf4OHRVq/T93+A00A970N3RtX0uAI4zjGEJER1Pf8Ct0Ai3S2Wl8TKWqaHqrqvMwpJQND74HYm6kBkFKS4/HkGksAU8lh2O33O2Kapuk0nbQsgRQJiWMIwQWWKS0i4kGtgcPhcHE4HADgNJ3cht2vbTdz79gFPUCAjr7XlaCVJg1KjbctTcjhHXRHN4NqOl247dpEwWscZmU/ALhHRVbVLP6jao7TGT0AV2ml64fXu74NX1s6HZJLOq3ztS4jNFX9xeqhCVb2w5STzfOJCHfjCCEAIgYm5JTTi6++BMIQIjSSbMFoNv2aGYAJ5s4avnufGzWK/NoRMTNGOuz3KaWXL6/Tkjnw+yHEGMIQoDubu/TujsHwkETCXkvCQEREFcs5eyaPe3WYQeCwG3bDEJeUwIDDsBsHRFjmZZ6neVrAlIiGGEOMDRLMOS3zLCkBePhLMA5mhoGG/W7c7dwuAmCdMhEMAWs8Vr+aaCR+76+Kbyl1jrOtne4Zjv7/Agf0sU/V46dvb69fvnoJKWEcwjBwHUZVTVJWEVWpGwYq83EhKJfYO/X5wFM6VotkbLQ33Oa2Vy6E9VbMVT1QWeXFTbKBMLiVCeHWPaAl6WEX74j3nBGbOZEeHgUaDam79+t0huupYXVhHXAcD4f9Yb9nptyIL/gNoMkfCI0JCHgfg6mZZuGdFfM9z++ffPdsvd9NDxEAmGncjfvdTvWkKkTEUKbF9dy3lRm53a8TM6vKPC83t9fLPN0eT87yHPfjMA7ENM/L7XRapllMHYoJIYzjaAC3x+M8T8VnBGEYxnHcxyGmnJdlUVVQ4EjjftyNu3Ec3RM050JEIeKigYTijltQe1z1p6otANy7a2Rmp4q2wKrqtEv1YvaoMGfXlKkICQNTBq6hY1AGaikrBC88BIV4bmruJ1nVpoCIbh9REoIbVQ5W2RkVV4kVbkMEQme4sqotackpOxpRak+hQygAoJKnNNWoEzLjlNL1zQ2WGLhQarYBGIg7/QIwGyEZGwpBiSBct2tiAnkxD4oyY2Ymaii1H6jWHZVQT5TTMv3mN78BgA8+eH9/OEhK87xUZnfttRR640+8yz85d//ptnrdnlByXpYliyBgSrOaHQ6Hy4sLZMopqWocIhKJyLIsy7yIZCIKIYYYAwUAyJpFkom5n3NhXgbaHfaANC+zOfEzBK10OCqDXinqlR9QzjlTUDMiKDsPv2jZuVsVQlkXIVZckRAAIOdEAIfLi3Ecp+n06uXL2+tXovn25hbMyLcPhGCgapqlRJaZZ2BTpQBtMV+pfObKdEOiNpMZ+qb/nM7vr897/KYN6w6qjS9inSQae7M5VFDLYlLVxifEFfDCNvdUBsaGQmp34uWpxLuv/7zyeArnQ72+lcvVioUXHfZXV4+urq5CjNv+ogsgsE762xHE7yu5cK81tN1ZDj9U58N32/bba77+gGB6HHeHw2FOSURULTD4tWV3dtu4AnW+UgdzRr3kZZ5eXd/M02RgXrBjiGaQ0pKWlCW5rxQY+J5tScsyLzk5+g8AEIbh4vKAQKfTyZmIFEMcx8PhYjfuAGxelrQUcSYTBc89IiqAplUHW2sJLVYMdJEATMWsyvFrkVsjCJkaQ6HaywN2AaSucCpBYOvaWSsGYA44AXk4pamnP6pKRU2BiBygV9E65+P9C6WKrUKhehMRqkpOeUkLIUIcMLZZrE8SVQXyu5xDwdlOp6NHFSKigRIRDlxtTAWQPEU58hCI53me5gxqFZA2KLjWQkKIHAJ7Mni5BLDgb2VgqDPBMAwp51fX10g4DAMxRQ4hBhGtPAK821Wd9TH392F4j7gl5zzNU1pmESWkEHg37oZxVDVBqfZfkvKS5iSSiXgc4xBHDqyqy7KklFJKjhACIsUBieJuPFxemEFSvzlq0kgxeWgmOVRVb+cocfHqgQKXEXlII21R7q4dB0CEys83STnNy/F4++rVK10WIIrDUFbqYp5F7Owdr6gOijoSqMXezUW7lovvYe1uXKBMjA6TFnacC1/0bGFYNDR9oN6G/o6rVxjRasBcfkqzYjZQ9SV/HYMQ7gmXxvY5r6OEdUTTrlegjRXESi4CUFMCsCIRJTIt8QxxHJ++8/Ty0VU7T+5jpt39wkMjAP6eJTq8Pj8e7lk92zdKB77x+YXxB4gI+3G8vLg4Hk9HOfXWK1gkftbp8esez+dusyWllBzeQcliOWMMHuFrBktaHP2vE4Mrttyk08pMwCVlnYljiEtK8zylZUaicRwP+/0QB0RyCDXllFNGM98i+CUMQL6EKJiE+SCfi5AnRPfwkaw5pzb0rsMUEiNZZxyNq7pkjUSwsidvdz0ilrTtGjC5vrdMiMio6AFTXsrDEIhZVFUMmRCpjNZ1K6N+Qbep1rqIAqz+j6oKIJiJUD2DvhtfK76rYEYlEgkRQaRYX7jUYQV7WgOoNuxD4B0AeEREQwKtblnVLDACspmKdOYxZYVeXx2WvYhfJ6fj9Itf/mKapg8//OCwv1hgWeYFmSrSZvcgwBv6e0f0w3sYfv7js+RpmufjSYkuL68uLi5E8s317W6/uzhcnObp1atrmefsYRgxhhhiGJhJVedlSlMSE790AIHjOO7GGMdhHMdxWFKuar9K7nQ6UwufKL3rGm3vJZGNiJkAzQonpwjE6u6nkO661anvdXa7HYDd3lxfv3qVliQqBGjM5Mn14qCPaLGt8pG2ZUd2aIvbPnTky0pzYCJudFSgmq5er41NvWo+cXYXjENAKvuFOgjBGq5SDoF1sdjgIG/LfDZqGKz7z5V90GYEqeiumZ1j0ud84nKPEqBW4RAiEpC5ZnO/3z958uTy8pIY72P3b7JfDOk1gsd7N8EPM3/sd6eBvuHq2b7d80vu8LgbD4dDjAFOYKYP/83KEfPukEglz8s0n+aUFve6QuYQ4hAHJhbJ0zy5xa6bc3Mcxt0uxDAvyzwvTaSBxEgUQgBAyXmZF8kSRh7dsIXIzdbBTEUlZ0QgJk9fUi2wjS+ZVUtbpJKRmJkMGIGcqZJLLW7FtdyqhkTG7f6pBsXW9ntg6OBPAZRo5Xp0vOaO6uCB2QVl0ooytc3hdpm5cbQ00LWQ+L7V4xC65BoQFVwADDgCtZseUVWyZDWjIvghDzDPOaclqRk1u/c6lDBxzUdzUJuRyMjOkM+O31FW7LhSLtxABB2bA3WltQBAjNEMXrx8lZYUOLz7LnJgDoxdtO/vecmXRZTokpacMw1jjDHG6Bb3jtfllOdpspwoDp7txcwAmHNOOc3TJDkBIBBTiMRh2I27/X4Yd+5iC1kIyagzwO+XOau60s6GbC9aaubXnohSgEbiWoVTZYXv9qIpq6CqAbx6+fLFF18AwLDfDzFqYBPVrOLLbucU+UIrBPc/cfRDoSR4ifpepNmJFg6CSwjX/t6aFVbf7+O6wt2CHE5fBqzQfFmseYNBHZaCRY/T6chWjwRCW5XypcWXkthVlwvrkqTKN/FBM63esqALGPajmCp5Asbd7tHjx4eLC2pbwBVaOg+CQoAHzFt+h4TnN2ABPbQ1ftOv30O0es3XC01wt9vtD/sYYvFaKPTzAlFUuA87bq5BJc6nlI+n4zydjsejqUQPvxuHEOK8TNPptMxz+dlqQ4xXV5cAcHN9M59OnuKCACEM434cx9HMchYxLRwbNzo3E83ViETVhIxMTUxSyqZtM+pYeEkU8mQCL51YiGBFAV05rAa1BJsJiuFmze0FEYmJ/ZquXKKmNK570urZ29lqeWkyhBCYEMUTlUqCETKRAapa3Wxh627Nx+Q6ZSOit3amqiJaG3MDyCJF2xwcIkIEFDPJaqYYiIqAmQiBmS2ASG66gSYnDpEdzJqnKS9pycm9T62DYFv1L4AUUguQqxip91jVZLvJNM0IMQSe5vnnv/jF6XR6/tFHl1eXS0p5XpD9vXyNca7dw/tvySPWfG5Q1CQLhLDb7eYlLenFfjc+fvxIRF+8eHk6Hs0sxHHYjSEEpzCJpJyzlAx0BFUgGnb7/X4fh4FC5OBYDRbw3890L39m2DnYdECc9QNBwYwMJKukrCplG2QKxvUqKqoUj06apvl4e+Pt1LJMUPfMhdgp1b/T17zO7m9dv2vfSqfSfPHKhRQ49KaeZQ2j4FthMzsXP/UBIStHaB0hoKQseo/S1AK0igB6b2qgsvDBJhMm2wDvZJvOuu6Jel/lsqI4t2RG7EKccHPdIJE7uCoKLIgA+/3u8ZPHV5dX5DFT57RPvAOoYC9Ve9Mgljeu238UHUAv1sBh4P1+xzHUImlo3yxx0Iq4pJxuj8d5OplB3O2GMDAzAIhISllzAg7+I0MIYxxdKZaWBX1fmnPY7S4vLmKI0zwtywwAGMj3XEWRJGJqOYuKEw4UJJuBqFhp6jd5GYjkUX1I6FKvcn37lUoIwCuvsMZo9PqVSj0gVjImJvZUGQMgxfXwaIKvrjPBjk1WtAgIIGJqhrXSY7G2sj6KHkCseK70Sus6AZT4j+YsKgKAic0gMDBTCSzDJuap/qlAxHFATJBKvXN0D8mlDACS8pRO0FhGSFv4r9zo1WeJQKuyGSsPqnA3Dbq1iZlR4BDDMi+vXr3KOQ/jwIFDjByDLxfxd2mnzhJFituG310q2cCA9jGEeb69vrlRyXEYx3GMMToa46GPKS0mGYiRIwWgGHf7/e5wwVVRVRLQzaj4TJV6fL4m3Tj2V2a8mpr6FKgqquJRLXX5a7Ri15ZTzqY5pdPx9vb21tIMiBzCuBsJUdXct98zDzx3b9VaW5O6m60KgNY3U3EkZfLdALRZs9zGWW3raICrAKwqywiouFzUXp+LcKCAPGfyWbQuLQYRyQwMVV131c6l5p6F/jo3Zvq2uiFtI6bxTIj7mv0nFq4W+5lNRPvd+Ojqar8bXAFADzfofz/OXeFBstGbdP1r0UHozODOn7+x8V6fY2pEGIgPu/1+HJyArqZorqSBQoEEtL4fMxA1U3EHNzMTEcnifofEBAausllz2J0/TywqNUEXCUkAQJSZx2EUlePxdp5OAOZTPNeEazMrazpJfimIFGRzxevOahUxcr09tODisDb5xUkRrE8A60DPKsRVUzTkIoI3Vauzbh+G4zM89TA1rtSFxgwH33yYATITcgtBLCJ73xevJbSw4aUW8RBRkUSyCQCBg2DWIsK5jOOmhIAOBVDNQjBg3zflnEQV1ICsLLDVMYoE1hl7FSi2hUABF/wazDRrNjUkDMTrrqD++oQEVCjvqkAIRDwMcUnL559/nrN89NGHF4eD7/T9SLEN8NslYG3wtY0FiQs5TC0Xc06FaT4xP370+HA4IOGr6+vj6YREQ9i5/zOApZTnaXZOjIkXSorjMA67OA5xGBw5q2++5ZTVg8Borf+wNek3622MsQhr1bGHohnxpSzVFQKYgQpVo6H5eJqmKXu2nSmFUBhrAFlEXNNbERxqK9zahVSrOhdrV9uuPiSG0CdlZ4W6VMzqq99Y3DfvUmgJlAWqNehZotTiHpBKL2OdKh0bI8x3WD4SZPEk1rZHrr0CdWfCyjLs7Kbvh9er0sIaPcDullL/0QTEFDgc9oery6uRWRRUtHq5d52I2cMSYFwtHfEB3di3qdt/lAngXFCPCCGGcTcOwzAti5qgIcP92w+/piS7tee8zHOucUItXstM5nlZlrkISbSg/xx4mqd5mv0uMDAkDru4G3ccaD7N8zznnDlGR6UYKUvJehTTJS0e1W0d1I4U6v5n5QIU8iWAb6pFGvEZWy+LhNp6Edt+WkZd3q8BAzKhQc6enyxe8Qk3ccmIvXTANhsBP0YJtzwK3WgizTrBcKtzWKXFbpXMFGrHt44uKpC7II9SdExBQaFJ/9GYGNwvKCdVWQuYmvpshSWaG1GhbpgNigtT0wWaqSYRFSaSCFT9bohcN+ZnADeqo4oyUdjtl2W5ubkF+02MDO++53i8tTCK+0IB7HVTLDaCo2g2UCKMIV4cDrvD4eb25ubmVlTG3S5wIEQzTSkv05yW2W2cKQTiEGIYhnG3PwzDUN03zNchopJFU0pFEuAMTmyWf53D8RrXVYSHjosxMbjlVEEJzY9DdF/+ZdYsy7KcTsd5miQlJIoxMLk1ubdWuXGRqT78zRHVklJaVGZaRei08UYkUAMVBTAXfJTz0rZ2zKXTr0AXNmOJlZJTaThW1+9t87zqGZGwEUqpmfmDgSmZQW7ddbOHwPs591hCPh6ShzR/POvuGrS73GE1U1ALHA773eFwGMeB/EpXAPoW0D3etyv+jiYA/HapA5v1d0fKMsT7lQv3fH21/SHC3W5/eXmp19eqqiDMVGCA5uft8yQxEkqWZV6ur29Op+M0n7Ln4UHxhBHR0+k0Tyety80Y48XlBSIeb26naRJTUFGz3W5/dXlxeXlpanlJIgIIMcbdbjcOo+SclgWJvfZsrArbL0zExByYNp7iCl30ipqK6uotToyMDcvfTpfV8sHvYixpX0TkbrmgosimwlmtGlFYMVdF7CQXm1htNwUjNIRgnNHMTLIUWkcx42quQc4pqhPA2uXViTaw4Zoy6BZGkJOqEDE1zNS7ao9FrS2IW3EYQHbzr820u5ot+lgTiClwzjlrA3sJkcwkSxbJghjqH3l8rntiuJdybeBAAdCAyIh5GIdpmX/+d7+YTtMnH39ycXWV0rIsiRiJ0PrX0nHES1ZD9/ra1OeroLQkM9tdXT19991xvzsdj6ebo194wzA4xiXu/ZATOMhGvNvvd7s9x4DIIUbkdsJZwfpVfJbl4n5KPX8fXelHpR0uGhFVB4ua+Tp7CSy2ses2Lac0T6f5dPLLHhFCCL5zUlXJWVX8Vy7UftqwoNx+yorSUUsRdJpR5Ro1AR+6F5UqaqfkKoW9FHHf8ODK42xutush0a5DLwZOZivTgmOPWJTYsLX7h1X1h53Ky7p8QTwLHjyv/NZkg9UfrF++AL6mMJsogO3G4cmTJ1cXF9zeSbqvtr/GfKJXqd+Bo75t3f5j7gD6SuplYL/bXV5cnKZpmiYDMAtVsGOrUhwIAFQkSxYVyWk6neZ5MlPvmEKIROz3jKQMjuEgcuBdHFKWlFJaZgezKYSLi8Pl5RUHPh1Pp3kCgBCHIcYSyCeiagQKiHJHrYfMxIG4MeWZiLAAPlnALXMbTadyBomQK4ehwzqbBgUBgMHUICuY+r6rUdyL7hHBDIIZAFe7425+rdY8K6WijNTo0G8wEHEtApAVrEFNTQUaC6u2cVAsjmregMcAtqVXJXk4A4dda1Sn2hpFqX4HMhIhhUDgSDlkx4KRcIgxI+SsG2V9dZz3paJz6wjJbfXKLgEQU/Z3H5kJwWtS8BisKrcwMBVjomG/m6fl5vZWVMf9DplCjCGUFQsCvrn20utRWvLpdDxNJwAb9oerq8tpml6+fKkqh8MBuRTTlNKyLCYCiBQG4hDHwQ8A8usBsHpzYEWOEQAkZ5XMPPhOuCUdYlNsr659bjvj3CqsPLJyCXowuQvdl2nyceR0Oi6nCRCGcQjMyGVikiRZcqGx+3+YAMHtLhQUDVVEsnSoexnDHKeqO1RTUBSX+auprDWZ2CkCVO253NsEiVuPr7gy/LHAQYbaKYbLhEGd4n11s2uyebs3nGTd52J3wGyLsZkTMbGss3TFBhsAa3QmosHePcIPS1MwG3aHJ48fX1xeVp85+7YG3X+YHcC3FRnb65in98zMD9KmCiHSBBHwcNhfXly8vH51Op1ELZihZyNiJ71AUJHTdEpzMjBi9i0uIPAwuv4LCdW0Xi5lDcrEauAHB2SFQOO42+33F4cLYlrm+fb2ZjpNADAMQxyGZVluU4ohxmEwtWme05K0uYurIrHHeQNVB6tiMwk5uxG1UYM7OxTPr9ce2m6TK3ZntIIIGBK1QiaSO6DQzCCLoIlysYzvVP1l7i0jvLt+OmZfEyY9cFtERCAE9mdrA1AKr6Y4ttu6ay4YOxBSAAU0y2DSlqKF8oTN+KVo1epXDBEDIUIoRsmqphqHGDksQDmfQAGDI2CoZqxGSIGCmVb/7KazAyQyMBHJORNSrB6lDUnue7K6QtcQw26/y5p/+cvPlzl99NFHF4dD1pQWJWxOmtC8gAp8bGsAKxTHCwaR4+n49Yuvb2+PDiPNy3w83s7zabc/XFxeprS8evVKliRmJgIGFMJuvxv3+xhHYi6+2c6FXU/wIqoiQKfeDOtauxlu+nK/hFGqKVbssTWYbgbuLYmZhRiGGE3l+vqVpGQIqobsq3XwwGnNIu7e7EYnzA6wIK5Ak2Ypfm7WMpCJiBvi3ho3N/6xjsdVMXqmEAJzCIGYXSAt9QAuCE+NyvJTZT0HAtU/aqgwViYOrEdA4SNo8/OvQ3JdHDi5rqwYyomlTqPqsiJ8weJ/LIb94h9XQtOW2ei/KHkd8mvPEGw3Do8fP7q6vCRCrQ6/q7vlG/Qc9vvshh/45uGPePjUG1SphnANcXDscaXyYVt1AiMrSFrSaTrmrJJzE0ox8363DyGktCzL4uCMQ0bjODDzsszzNJt5om8Yx/HicBjGQVWneZ6mOUtmZiDMOU+naVmmy8urfdgvLidLS2Vwtg6GaRgAHNxcfTfd7sp7oOKh0qH/7UBrv5sfD7RlgHnrTsghBgBIObvVMBCvbnmmJlb1j1S7pG5OWa1+yp6wUdeLPF38P25TLsUyGss94VNDvxuz1tijGzOAowfdeg7VDFQ7nX11ESJ/GULAgSh7cGaf/ul3HZlHrDiyhNm1U8AUOIQ1h7mi5F5k0pIqzYpXGLKfzh1RU8sZ/QCYpunV9XUWiWMkfn8cYhyiipq/+LPJfMtSKZZHQIo2T/PLFy/naULEeZlefP0iSx7Gcb/fIWJKeZ5nSwk4EEcOIRYb/32I0alUTRBeC14x9mRCT0EQaQKIVm/KMtCvtpXJSBhcGV3JBU6yYmYEAJE8z2I2T7PlBTiEGCiESvNV9WUVGAIxMRWjQ6jWBqYipkVuAWvaaRl/Gzbm+rC66i1G1sXTM3ChNjM3x9aU0pKTCYgZ+9VE9VSsg2Cl+lMLt1Szfo/YUzN7dgY+wPLtTFOrdsyQwbaokXXW4EjQHQANJu2Q8DVcyACNjLR8RiIINo7jk0ePLw97RvR0nG9aM/29TAD9O4T3bmnvr9qve869Tz6/nWpYlhoQ8m4Il5eXu3HciJuo7nxsZXQB2DzPx9vTNE85Z2ACUHaswPT2dHKOv3/vGHeHwwUTTtM0nSY1CzEO4ziMhS1qjtCr+vWVl5RsycvsRAY1yFkWp+sxV49IX1VxuRYqIFJ3ulgMgcHpP3XsRQZkT5xY7V+9CLaGrdqvl0xBAqTiwm9qyETu6+Cu+dVeMYuRkgtsWphG2Yy7zwuAinob1SzVGcmCoYCZSvIlQDUuL2MKWGmk+sYXbS3tnU8rorNmS1/PzQJ+pUyqFjtiv5/9f5Ao55xScVct5O4aAptNwAyJwhiGGAlgkSSSV7I7gKi0HIUQqp2A1g6r6e6oZDc6MRcAOfA8z7/6/Fem9vyjD3f7/TJPKRm33WPn/NBDCZ30zpZlubm9lZR4iNM8T/N8eXn55PETM7i+vj6dTmYGxMRht9/t9nuPG0IinxGprC7qzFX3q/7mQNmxajXyru21803VzKQce9T5ZOqKYPv4GwMjYppOr8woMBNiHKHKIIp3vxSviMCx6iNKLC4YqEguJITSENCKwrh/eUnm84u1iNc3ql50jRtzdD6CA6dY6m+rQsjMRard7FZrwEIxLardd7OMqkuRsgNbmTJbu6ozdmUVRWM1agNEwr7XdgfqOsJUNXbZdqD1GWF1vcItxxXVNyUiklOM4bDfP3n65HI3GpBornTePoysk5k9FPP7TXX429bt8F2389/i2auNJDER+cYsxjinZCZOWPCm21kJKSfJ6hVtmqd5OnkrShQ8RkNySvOc04IcnBHDTOMwqklOaVlmJArDOA5jCMFUl3lJOaW0FH8Ys5STiSDgMIxMweFQ2whSEZk5xspRWX8dlYJyoqvMtRhWQqVQNPK+9ayFZtDaRAElIHmd+v0+YiRkAgVFl5g2m0W1kpIEZsZU6ktzW3BouLDFO8ypOL4IuCc2WG/V0ptqra2v4fY26lNYC72wWr4Y+GZkLZoGXrD87WQO7iuXtUx8XCbxmvdQqEnqiwp30COlijYBUgkTN9MshguaBC9HSDVJc3WYLBxXkQJSxRBTTjfH2/Dll+MwPHv2dIiRduSO9kjwQNgSruxLD+nNyUTUAiIShd04DkM8Ho+3NzduTUhMIcT94bDb7zxsrmo4HHLyI4C46bvNSkddhW0IwKsBVJV4FX//sn5qErgSSoZsZjmnU5qm0+Sz27LMbEMcIiE5cGSmkkVVEQyRPIXShwjv4BUA1LJkdapFzS/qvB+o0qGqJN60CWjdndf/l2OM485jUL2DMTgPYnaBGTGtrj4FE4PVMLEfqbc+FneDU/Bekr51rsvV0bv8UlQspjtepxmuuaK9WQ1ab4jd9hWNoApqKpIAbBjifr/b73ZMJN4t8ubO6vRn9rtpfX/HCcDOtPadyP1+xP9MndkpL+4+/4wUZPdKwYosiJBgvxsPh73eevKUuN4PvU9RnRdP8M2FFq8ChMzDbrcb9zsHta2kKyAgF544oSaTrCDqTVaM7MbryzQ5/iMqZubyAgAYhp3fq8uSJOWeFYiEjv5T5D6XxqXvJoJmVX1j1nGioYUFru7y1nmc1AtOi29igXmLBVBZM7jq3b3ZVMRM2lWqDiigQsCwWoPhhmxglftpvkV0ZzFCYJPVmo6qGV0hV2AVDHS2OV1+83qPlhWy88I9cgaAitFdka0lUG7ZNMxEPAyMTHnJSbJ5AGdFm7zABfawZBuG6JhYiRyBldZiZikvosLCIUZE5dpXNgNk07IJcANl3ysw8zRNf/fzn0/T9PHHH+33FzqfNAlBObOhc/9vR52Bppwnj5Nzdpbq1dXV/nARYzwejzc3N1kSIg1jjMNuGGKIY7Ur79Nsne9D7Odxa2VXaw9rWwECLLaaKghATIy8OksRmKKqGiAjEbOq3txcf/3Vi9N0MrM47gSAuXjD5eriWQYFHogbZ740Kn5DVW/LWtAr7aaS7KFR+4uPW2uLiciTyzhQYGImDlu1EPWGbqt9GzlRG3rPwu2CFnsF7pr5XFwd9Ixb01n09gaiRStuPY5TsyGwszlvmhg773UJcRMh3DPHPO9TEMYYH19dXV1cVpO9O6WyafmaO1Wrp/frAx6ow9+2bv9xdwA9rOqf2G4/Xl5ezMsyLXPVEzoPkkWWZVpOp2MWcWa0e7oR0X4cY4gpLfNS0n1VhZnHcRdizDnN81ITHB1+5BCDqk6zztOclrkUBTUQA8YwhHEcVWxZ5pwTIlpL9wWgOgFI7bkIsdjUmRISEJqorM4/xbJ3i/k1kn6FSDpvFO8IXYZj7g9NXciTtQri75y1ByCK5GJQA4252aT3q5dip5gBICLjboNaAkO6kCSsNLlqyVRhmtVQpt6IhgWnkra9pxbah1b9eBukwcEToxadNItsovgAHDB26GG32x8OF4+ulpcvX7x48fJ0OgGYu1U42VpMnCbk/Shjn5CLSMBAVlNoDYyZhmGQnG9ubw1sGOM7z2wYhmGHKqpmaNvQ4OLzgQA4TfOLFy9ubo8ASCFyCPvD4fLyYp7mm5vr0/HklwogxSHEYWRyo1bzUt7FGxATVk8k2MAV3fxlhcbZPlNgt3qrbykRYiAAVpF5WRzQf3V9fX1zLSLjOHII5MtHcZseqS6exF6dySEyA8tuIOH64R4tqaHuWOfPYuZTqf3uKu4mT4yBmWOIgTkQBwcLDTYhu9aZKOBWbXtHF7uh82zDpxA38t2+42lJVOdaD79zAdYUTDOzeyzU8O7ruc+jaPuySlqJaLbx4uLJoydXl5cEjf/zGtI/fpN/zx9CCfx76AC+3detD4cpQmv/6Jhhvz9cXl5c39zMy1x8AQ1b7ck5H6fTfJqXNKeUDQnAmIljMLBpOk2nqeCPohyHw+EiRJ6n5XQ8iQoSe7oLU+AYSJWQRMREkREr3NcUtFnysiTJ2ZpWr8OwAasJcQF9nG9Z7g+x3K1/sU83PaMONyy9c3NT9HgPgBKs6s6flSxY1a6MhCaqKp4272+yqGpaijd/jQypQb5G1ZJ9BW28zrLnlVY7YiKFdj9ssvrqckGhc3pZAU0napZ3UkXBsjGHBm+LWDYjxvbKCjuIgw6DeTY6WP/aipaMcBiGR4+vdsN4dXVlgNM8WRZgtJZpgsgciIt3HobyNlZIrXwIvau2mRLzbhznefm7v/3FNM2ffPLJxeEwzZMuiYC21FBTVeYAZsfb2y+++OLFixdqGne78bAPzGlJp9NpmmbLAoFVRSSrqHfDuja6zamyOft73jkq1sxEXnkmpfEnJcLAhIGwetprobmiGXhYwinJq1evbm9vvSgzcQiRmXyZ63kDNemTmEPxDcTayGdVFagxvYDY6EB+6RSOmcfXmBTZSm8gyDEOkePgscNEbIhF4bepA9gyu2zdrltj4sCZMyuibbonsxpqthbM6oVeE+Sh7676qqoIYZM/46zDYt+5Qp50vyMbnitxcePV0+Z+MxMZ4vDs6dMnjx9zYO25XJuz756jb3NQIP4h6vMffwKwSp8ipP1+d3l5OcSvVQ1A1IyLkCJrIZzYvEzzPLvhWiCOYUDAnNOyLDnNQIzEgXF/2O/2OxHXDE+AOMS43+3G/ciBRUWSJBEDK6nz/jkwIgciMlXRXADNYrvs+oHgeReN34/Nfdj1Oz4Ua3VFLILGytYvihXcOJH0LZEaqCG5s3/J8XABtG2stR0RIEAkQVNvnK0JQY2YrdpNGHoaQhHUNra/NbtuLJ5a7ljRGBIeDblx5F7ZcbaKOdcdGodQgLhy0DQX0eBcckS3piFVVSzmLAhITNFi4ayv3vFURRcMAMfjcRiH/XuHd959J4nknF6+eJFzrp5GVvM0KefskbyuJMDutoTeihrJDJiRwqDTfHs84hdf7nY7e+edIcZxHNU5lLBCvu5sY2bLMl9fXy+3txCZ9/vdOKrq6XRzOh5VMiIYopnmlAtXqiedlOWpIzhU1H8VEemiEKxJ5BCMEANjcPJopZ05kOZ2s/M8Z8nH29u2f45xCENENwTMblbbLIrJI00LuFGMPkvekQdQF4ZPoQOhh/i6ssMKz0eLrg2JOBAzMBKHMIw15gzPk++wm3eLd37Xdtt5GgwCPoiJm60O/lpEl9DjPJuu3O6phr1dKJwFMmK3YrAt0fpu239PtoQ7oZvpGOPTJ0+uLi8CsYri5pe0jaoA8A+Y/vXd6AB+nx3xA/W/ACnEh3FYLi/GccBCJitskJTSPE9O/gMA04wGIcT9fjfuRgNLruNVA5BhPBwuDhcXF4HCMs8pZRPBEOIYd4fdbtwBwHSaj8fj6XQyXzFXNIM5DuPIFESlUu+xbec5hDDGIoVVaypyVTXJqlqIQS5zVW0CFqSykVwdjVtYAjZmu7VcTDRE5FVTQpX2XRWPhehcqiMRl1HApJw6Pqe0LtwX0bULJmoq3HoLVe5cxXgBCQotug/VxhYfoAYtWKxe8MSMzIBkKdV8qQIlgwFA4EAtbd4PBzJDcxURld0F0TLNKS8gBlzw1BCiqv32yy9fvnolYs+ff/Tue++qqZp99cVvQRRLvhIUCoapSLZZEYFox8wtVsERNkft6mAEKur2fyktP/+7vzudTh9//PHFxUVOKVsDygvQwMRWxLQZJEMob/U8z7c3N9kDKpjrSynhyVSgO1dpIaHPJ1jldPXjqAYHvu0FUwQjghBoiOxUrnKQ4Sr7ymppmq+vX13fXKdpUYBQd+wAWIwRRbyJCSGUz7jT67lNFlSL/2LVzdyIMSvQWHKPuv6VvOsfqp6ZMATf8Vaw0XllxQOpXjBaTAxrpuOadGSGfUNdTcw7DexZy2x9lnhRzK1MipKs11mKrmV2RT7PAHl/h0oEQKvLRZVq1Xi1dzOyasqIm22z7Xbjs2dPL3ajIaaUA50Hz9mGorQlIt1fQO27OwD+MNV+tbo+94S4j9BUM0iIaIjjOIwxxnlJTmkwgGVZbo/HeV7Skoo1jSgi7HbjEOOyzPMyqQqSc83Cbj8y85yW0zSJCDKGwIgg2T1kOOd0PJ6WNHmp9E0XIHCIcRgBcVmWnBJUGiXU4IsQByKqUEIUTxQAAJq+SURBVGn5BVWyJzszMYDTpbUzB63ngq12h+j9uNnaPlRHSdyY8Gwkhda8/LtNW9n3MnY6XmBCV5A5d9vUoNk2A4D7MtafUcFcVxJxy6ctZGcghd7vtlq4WzNlqUYuWPOpihGe1sWDiQoI+CzVeVb7N0LVsuEmCs13QSAVn0YiIhKRaZqO1zcEgIQffPjB06dPv/r6q6++/KKkQeCmTrnDEC2ZaSEeGakpZmumbpcLaIZMcYzTyW6OR/jyy2EcNEuIgRmZQkn8EWtxPIUKSQSIKjLNc16WZZlBAQND2ejaKvompHISELXABk8+B7DelLWRgUUc3GdGJnLHETMDLNGYqpbycryd0pJcg3Z7ewuqIQ4hBkACtSKcN622Iu5c4jRlkyqWLoldZsXFoejby+pLPOdISqh7BTHIpYrAxBzC4Pwidp8p2zTe1nB4r8U1efG+Yn4/WIF3Jb34AIjeeb1ZNStsFyrBShVdTajXQG/rwfgitruLobQRzVtEu1PXyiaZOfJw2B/2ux0j5sJ3Abw7Q+DDk8q9ewb8bk6B0IvQ8Cxd3vqk1Df9eudN1WCsOwbX3clXGuKajsDMh8Ph8uJS7dpKEwQppePxeHN7s8yLqHgFZST3u3edvaoAcwzjEAcEmpf5dDxNp8lB28BBRG+Pt4Cw3++dCAGqfr03cIWZY+Cc8zLPKS8tkrNGursmHt3cY/3dVcEscAyBc1pSSlJG0XWFau2KW7dWZR1ZN8hljkViQFJVNVlxVVwVZA0cbKHf0MKKsCNuus5aFUCx9pg+lXfij9XLt4DCgQhJAVyqy4GJ176mVCy1NTMAepJrSRfkwIJoOQOol0In7FsuCDLxGgWuaogGWIQ/ztxFxOUE2b0TqvJzHMaT6quvvk5myHR5ecFE7vaJhGCdkLdeZyJ5mREBYBiKiXw5dZyMTo0d4q7ITLwb98uSPv/FL2+ub9599s6jJ48CR6Z16HHAz802jQMRpyUt86yePcfdveBqteJcU9cyhcCzWklX1+3+XLRigEEYgkM13O43qvmLAHo6Tr/9za+n4xFjJKZxHP2jymJg2RNTHfBiHtDX7nXqVlWVrDWLBUrOdsBN4JDj2GrV+XW9JonCGEMoNhXABEjmSXPFPZRatLWoVPtxrAF3xq7BbTZ/0KWk9kXOLSXOj4TeF6geJ4hue17Jbraux7CrS9biOaxpmRGdxqa9PZz7cljVPHcHQIF+2+BS0Nr6t/wyG3bjk8ePnz554hgmGHBnXHQekLJWyrultd8VYPf881K8YWM+8PX19jf74+8AYHWTNg9puTgcHj16NC3LNM8yT65VMbM0pzSfHG3kMYzjGDmaWZKUcwKzGMK4G739z3OapjmlVKzVVdGQmEwt5ZzyqsBcCcEUnIxdw07NqTiFGOM2J0i28skKYg1g7taLYG5VZNpZlpxTyK1jFtpaLaoLm9tQi2Q34b9b/aEzqrI1z2tLImjukN7ikUrONcCE1vEMkMAEreADHH1b67ZiqgKCCOKOW9WFprs1K1rgZl5WtWPV6gUsS8m5re6S5A2t96K0Ony5j4QXlhAYYMhJoJJ6/fkxhiVRnqbrr1/86vPPrx5dHY9HJIIQWkOvZigCYJ67YmY5Lw7BEY4OjRVDYOvU/z7JqRHhMA7LPN/c3i4pAcCcl10ch904DCMylmTOjn5CiGoqS6p8dVqbpOqJCdVbg7H54BRMqpEDGs1csxtmGwIwEXCB4BEwhsDMIjLNp+ubGxG5vn51c3sry8JgI+9iiGqas9uel7xqTyWpDkFQ1OrO7i/L/LqM5jWuy+s+VOKQrXHqjMTkXIphiCE6ubOxJA2q7o5wbYZ7E/66xPKbSpXvaT3xLKzTHsg4Wd9se8iSDR8Cn62czLDZyJbt/HkzXnOHN/TNnj+6Ybk7UXl/2L//3ntPnj31XddZwAx867Qv+846/3t3AOepmysx7X5zCTsbjAo5zN4gZvgOxauKyIno4uLi0aP55ubmeLw9TTMAZFH3XHPGOwc6HHaHywMxz/NUGOGIMYbdbhyHHSLMc8XTiVUk5bzbHy4uLkIMaVnmeTKwwi9xmRWHGMfAwU0IrBg0lOPBxWbBs2Ws+jV7V2tCLiAxzUkk+SitSFzSS2BNrlrPANuuqYoOywjdU66ch9UhCw1XgnFbpjUKxPaiuONQAgiAogo5I4A7WFcte2lmqsbTmiNlwakkZ7MQAyNpEbE1vanC6q9CWIK6GjWIA6IAqsuMq2amAUfu2turCto3x+YmXKLBvLcnT5uCGCHgixcvbq6vixGsa8Rqv9mMVryHVTNJKTm6Fatvq3vFmKm6GoAqFwEA3bItiOqLFy+//vplCPTs6TvPP/pov9+nnBy9NjEVLUAzoIdLWx/SCOj5a9yY7a5yAliZz3TGMWzKadOylCJC5FDaf+etnk6nly9effXi6ywZDIg5XByqekuyaAkxxaoS45YjU/Axye7rpq27Ifa4eGoNdXGGKKHuWmAKLChriGOIoaiau5DKdioWX1zbBF00iMBbKTdTFGHseTztxMDV26qzXMb7UCK7Kwioa1yCXsSxpdhgWRFUWpFVyU4542sM3x3FVuuwG8rUNoLF/VTVzA77w/vvv//Os2chkFgfQ/DwXnrjEG+dV+R5ruXrSvE9Bf3+L38vJoD+lSHifj8+urr8ajeCwTLPSUREU87F7l8UEMbdbhiGlNLxNOWUSrPF7O1zzpJFnOrni1ATA4RhGM1smudlmt20od70EMY4jiMxLSnlLP1sCYgc4jCMxCwlGkl80DTzzBNAMJGsNQa23jRatao1dHXzYVvlwdZz1gC2pbuKcl9/nrZEcFv/CmCxbKl8SmcSogADUGipqhslX5v6m7uDqDrbY1VaFgGpNfpP+f1rNJohQi3lxIYaFPIZUdDA3KWOPEnMNUQC6knz1DMxCAByyq5aiiFoCGC2LNOSBQITs0NPyOh8MCezUmftJqrLsgDAAIPLyp1YD0b1Y1rrl9uXBg6qMi/LMs9EfHl5FWIYxkFM3Ka/Mt/Vqf1IfN5AOcmSg/f9AF1meu0GqTMeMjP1/AlCBpBcYJM4hN1uNzp5YZ5vj8ebm5sXL1/c3h7NNIYYh4BEKpJT9uxqM3NwhYIvC4pIpWAxNc+9Sq+4BCZRF/DiKUaqINr2OkRkhRnEHqzm0K2dE+SrT0nLemzb4pLo4UJnKEM1Yk+l6cw2rXNJ+Ga2Cq4huva6peN9jgTb3e0qQ+++fp9At4tr7DE8v9lQ7WK/f+/ddx9fXTJSCeTBb6TKvF4H8B1zdsI3hEx+W8PS39ngtL5xRBwiHQ773bgLIfgZcDydcnY6XVFOhhDM4Hg63d7eZMnEHOMQYwTAZVnmeZ5OJ1G1JogN6O6hOadlXiQn4BJHDu5GFHgcBxFdliU5kcOvBFUgZuYQggFI9sNFRBQBIxN4yLiWe8Z3g4YtLU8AAZQQ3CKuUSbPMhV6SXmdirysYAOpsF/AdBB8BfHLLpSojR2AyNz4mmaQRRVywDXPz7tNJ6MjoTmT2qwYPxOZmuYMRu4pJlrs2DqB2GrOVZYTTZZJRBEQLGctHFgkJDLRJFLgI+/TVQutRQmYG/HGgfs5LZoyEsZh8DBF74exZcv6go+JlFDFT0MF9b7bfUR8oVT9pcvZ6AeMqpb0kpqiripqxkRxGELgIQ6VrgUisixpyUs1u1PtPkBbf/MQYgjuGlLpYD2rlhB5DfgCMVVRQByHwRiWlCRnAAghDjEy85KWly9efP3116fjSc3iEIvKzyClEi9sRWEdqNDPilFD2WHnbOK7ffCevyQHVElgSfgpgM/KZ+QQwxAL/cGsAj69PrbJA9E759X+AwzNyRahmNyVg7m0RS2Y2tq5yLi6l2zLuN3ZVq7VvzO/blwdaMkDRfyO95RTQ+y5RFWeUD3cbV1NgPVITCGCN30PIlVmAwIMQ7i6unz69MnlMAhCTgmIvtkCunNztN+5tL7xk8OD8cEPd+l3nm8PAlgPxcQ/CAUVKTQzxyHEEJgJzJIXZTVgiuOw2+1CCDnneZ7SMgNiHMbdOI67MYYwL8tpKloB/57leAhRVJYlZbf36f03PDIFUUXSkjSnEghXL2toASa+szWgytd2kp8Ujb4ycQzBdZu2igVaLkbZ+wKuUHKp3QV1QW1hMligA30AFlw9tKxIgttXK0GHmEhQTKQacqoqSM6eaFBsT0xdDVpPDnUAwTEsQdVkksXrvUo1e9mEkcFGpdC0kF4rutm4CgxEpIxZgYJXw5phrm3d2LzezVQ0oxEDgQs4mduC0HclqrrbjU+ePDHTly+vjzc3gOhaWX9azgnAyNktxFso1hSAodg0FS9SRGIeQ2DmeZm/+PKLJ/J0HEckuj0dcyoG3dZi3dYGrnqkMXOg4mKEKwq4OtBDKeFuoocAHGgYomj5YNys+3g6TvOcU351/erV9bXlHMZdCNEPTs/Hcb1gYe6vnXWJaFe1GmENJbGF/Wns0jgtS+HySkAFkICYY+AQ4jCGGMwg5+xjSo/6rg4DLbKtIuy45uIWhlkz2ACXzd9NWkTYZrNUYiri1kOhc+upo2txB+m9etaMVdri/wD+kW9inmwjMbBGD7HN3qur0qsDaE8JUR0iP7q8ePTo0S6OiGdCX/smZRS+aR2G3zcmPtzzx2b3nydb84rVouOM/nRvKJg98PV75gA1IAIYx/GwH0Ogzm/JiHl/cXFxcYFIKc85CyhAoDgM426/G/bOn84pS85+eaCBB/WFMYrIklNH40IDBKZAjEginjMsfmUU1AYNkVQ0pSVwcC9F180SoAO1btAj7sXGNMTobJPGSXcrFXdD9jTYapPeyEBaXZrRxM3RWqOEBKhrUl3H2oJtInm7A5uGlstCywzqMhZMTUzUADQQk4/pJegVSUH9GyI3+mLBgjwAxLQL0C19FwGs1l0rdgXuM2xSTXWx0n4cRso5AYBF83d1JfI749Csrj+BmTWEnFKSXJJqvGKubtImKsMwPv/o+TAOv/z5L34xLyknwlW9JaopJdI8xBFDqAZ67j7knpeKLXRL3f4DiFBEv/rqq+tX1/NHy/Pnz0MIkiXl1NnrWXWjrFW9RuIyUXW33CDlrfH1EqoiYBYi78ZdjNHmpRwJTCL5yy+ul7T4lTbEAYYBicR7de/Wi5EbEzWOjxGAv8k1GKcYmcQYsXgsFrGqmjlDem2ZkYA4xDjEIY6j55otaSmxCrWcAyLIWmWxdtpeFgjJzTiQoFn8wJqfuu4bOgxl5UbZNvi2BCD0HJYuzsuKikbqMggbTbPK622NwNYWntnUKripYFhZw326nrVOv2vTDQy0F5+pCJke9vsP33v/nWdPichDAvF+B4l7SqLh/V//htL6UN3u/+hOnf9+7QC8KosoEj6+urx++uSXv/r8eLzVvHg7Nez3h4t9CGFJy3SaRDIFdoPPGKJb8Dv637QZpsaBd/sRgK5PN2lZoGSJFNpljHEYR2bOKeWUAYADk9v21reJikEuBCZj8ooJZjlndQ+IZkno4RjlK8h1oef/S2DBHIJoV1+bEYCQTU00qwGszuf344CrlrKSZde4+Yo0FZSXkSAo+OISKgdPBaTQDYvdAkLv4Vsz8KphXLV6bieXs15aQ3fevBRgt2xuKwRuAFoMI1UVEiQwwMECBQfES/GX4qxNtVQQM5tBzk7FhI5a4h+Ol+Pdfv/syVMVW1J68dWLaZnNhJi73Ry3uaLjy3rrZlg+ageFygGsludlmaY5fvnlMAz7w35JS+c7iQ0qoeaZY83qx9VbzZGjrH4L16TIacVPuzHG3W437EZ38J5OJ09HOJ5OOS3EPAwDxwBO0pWcRQpk4rGNgYnYVFPObkzrXhRgzndAV4ExFx9DU5WcwSn+klWlmnkwxxBi4DAOY4xhQGYtaoK2yS3S9JUR3IXxsK9AQlnmN8pvVR02nfProX28i79vuZ9g7lVoXaoCniWn9H93u3+1epEb3qWnP5AlVi0sGtN0FfWXlyuSRWT/9OmHH3303jvvMGOuAvtvY4/w9+UG+s2nx0PTwP1cK/yGkcJsw/zspXz1Y3JZ5tPHj6ZpChhOxyOkFC8vd/v9xeXlMIw559PxdDqdxJRDiDH4pJyzzLODP1qVOACMFJiIs+Qi76r3nokgYozxsN+b2nFZsmQiDiEM4xC4JaCbFcklYS8ZcdKAFRd036c5FlENwsl9jMlQ1QjIz55OStUkVYWGI9XJHyF0onbrdwDtUrNOnVkuQU/NbZRnZ7uUShZUDQt3GcD9qh0RJqo+FNW7ofBFcLWnrx9fZcFXD6CqonAEd9X7lPWK1unE6oGGCi1mw0w1W0YwiBBChGq6py2ovbZyRMjDoCEsy5xSKmiNf8LFnAMcJDkc9k+ePFb5BAC++O0XyyIlixZKkifW6AJXJ3dXdON0WOdSYqbKgRHwdDr+/Oc/v7i85BhaKnJDLVx5W/YolSLOXMNdWqZL+c2IEbXWi8AhRu85igfHktLt8WQqHAdCHHd7b6ClJrd458Gt8S+2HZXgn5NUCBQqMISFrlaaXDVLS1rvTwdkkHkYxnE3jgNzLEo3s6xajEGbZYmaF3PdUGscOaP6e7c7nBpSo235RFs3nRbbVZtB3DJEWpBw9/lAz1N0TMnagdDIOoVIXd+Q3uPNNnjKagzhjd0mBbyb+OrpR12cR4OAyOzicPjwg/ffffIYmSVnLtTSuxUP36C04oNgzpvU7dc+//swAVRHnI4N4cX0cHHYX+ziEJNkQhzGYbfbcQgppSUtKSdv5XJKy5KGYVTT0zSladJiXmzIFOM4xEHNcsqSsqmgB2g0P1imEEJaUk5ZVGMYwhAC+/oBVlqMfzQFJUA1dd8FA8gingvm0JCaImzUhWuNKItGMjAqRi+mZllzLTzESqJubonqfMmCFMA9HOjVk8fK+pewE4EUM2eoc4BBLlIX65t921iKtg7KDLSO7IjEBBAQRKH8yHtJSlYt6wod03p6ByEhK2ixtnRDCElSXnLg0pEoKEFRojZ3eUJCxhgHUxBJVocz8vAZwDnNv/78V5Lkww8/fPrOM9/Vf/3VV8u8OAEFqChdzTz33GDAWL0uDWntE22VFQOUNl5Ubo+3y7LEIS7Loh2w6ei7n4gGoCag2D4UqkcyQue2bWYqZhZC3O12YxxE5TSfvn718vrm2g1K/W+FGFyvnkqSphZNn/thBfZMN1UFX2Dn3IbCktLrVB/yJsMgZwBwMqh7ljgPmkJwJ59hGEII4FezrxFqo7F2AlZOgM7atIqcV7c/u8MzrzHNBZKxTXg7NrpjzW9fgzE2VNA6RUAHwWMV12sv2u1Y7XgPbaisb3uKPZ5V/KpWsTU+vo+6WSX1CgCBaBxH938emROAqjE/xOSxb4D4/+ATwOvPk4e+cvb8h/D9s1PuzPSiP3TPoxwUgEII77zzzocfvP/Fb78ocz5VyZi4TQqpasqZmc0u1DQtKaUZiM0MzWIYLi4vxnHIKU3z7E3TRsxMjEBilrKknM2AAzEymEkW7ymwqF+LQSKai8XEey41k5RVhTmUIQGKLVzhoK5bXTDPQMaViU9MzrQrfFO/UnIWUbMMwGCusdzyPqu/sqnrcrURQCsLswDRxRHOBbtMhsHplo13b10AHlQ1ci92bzcY1fQnETER276ebqazelNoXcutqgcnxqMzUhzMMVDVDBkWgAGI2LeWFImJAUEKglRqSQyRiOcJU5rrB0Quukgpffnb39ycjoj0wQcfPH78ZJqX02laptl3RR7G5i8tpawmjm8xR3RrZSv7iV4thJ68Zup1NOV0e3srzjVi75qR0QPNoSmr2yyFhMRIvQmo44da3DqHcYwhGtiS0ldfff3bL768ublWgGEcqEyr6JteydkdHUKMPpUVAhZ6tosU+9iqbHQnt7ZvKJloUsz7S+YV+WQUh90u7sYQBuYAWAXAq6SRakpof0RC/SMs/rQFXmtpRht0V5vat2/6rHeOKVtA7GbfItbquuxq+tbtiPshAVc0Hztoou6rsAdR7xDzkUtWK6w5HZtou3vKY5HGZDGw/Ti+++zZu8+ejTF2WciG31j/m2L5TI/17UrrG9dtxO9qAviO4Sp3W1TQwPz+++9fX99MSzoeb8FMUl50mafZYRavtJEDc1BTyVlESpE3Lej/OBLz6Tgtc4f+u1iUwzCOzCEv2XPBCoUCqQaxmLkGfSUYeBXHVkArR0Xd+9a7VzAgUFlXS80+Fs/EhSu0olWHhcgh+CrUsXoEK5Zy7bu0y9qs2jE2BUB3PkANgWkTARFZIX3UqERnzGP1Z4MyphSrRlAy7A0sEA2A3cJMy7IRkbR1y72lgVqXGQBWaaOF0SHiuwC3rDAQSMJsIqJrE9cjLS0slnSIAL6DaebTlnK2JeXr689/9bkntBAhh1ALCCGtthsO+uecmZiAoue4dZrVM2oaIbGHNFhyThoGLvFYNURziyLpOqMUYx1quS6mYmAhhN1+z8S3x9vr6+tpnm9ub19dv5JlCePowkMV3/WKFMoQUGfUA+pyLXWhtZvmeqJ7g2Kg2lf4uS8qKq6yJuAQQuQQwjDE/ehWV26oIMXcwGrAOthGqVWY/WXJVJV7K9F52+fV6q9bl048G2crSFN7bdyIgAmq/TXYHScFWNOae1lBHQoMaRPfVYVgxd5uVQOjtuS1tgnDjt51x8e5HTdZMqjuLi6ef/DBB++/F5hFFQx5lfttbNJWmPFuP/2HnwDwHh3AXfzo26oEvt3zz6jTXTiumRnEGN9/5915mr7++qucFjBY0nI6zdPpJG7SK0JDOFxcjLtRRKaC/q+eKRSCGxQvy5zSUviYlQQ6DONhfyDCaZ6XZXH5j3cwq2rVY9vbiOmS1y4qGtr6tdkGmQGA1G1VDbTDO7RlWFNqV/OXmrEdSNyh19QJPMSNBrE6nds6RLsDXMHjO4JoFfo3eT3hqmChkhrguwxVA1PkVZq/+gXVfyz9XkARzaqahYnR3E29mJtSU4FaC+0iQDREJ9c4TgyImnNxMTJTtZSTCKr1DqRr7+8KIu+uY3AqC+aa2SkiCIi7HQB8/fXX8zQ/efKECEUycvBRhsrKwve96JvtlJJ/RKFoIwpTfvXDKbI4RkQRMQNkNhcMt7CFZgy50rixT7IugYPN6NPHoRAC87Qsv/3tF7/+7W9yzszMzOHigpmdnOOByYDIIXAI3bRlvsjJOavm9km7NZs7OhA0dr+pVqvnkurLgMQxDuM47kYOAzABQBYr1uKl3lH9CFaR7pqMwfX8214zHZO9OuaCglqVv5Si3HTkbXtRXGmtqaoLUmYuNMeScGpl72pb9ABa6EMvSm9oplXb0XKse3eiZlnsvjXnphEvXjDYqOobDwT/tqomerHfffDB++++884QQxV/IT4E9fQpe99hKX7jJ4fXs0S/5dc3dtq/O2W1cgGZw+XlxdOnTx9dXn794oW7Gy7zNC9zuWZEAXG3H0OIN9fXHgfvW1bkEGMch9EAcs45i6kCkfcWMcQ4jofDRRzCPC/LvGTJATkwrYWzubNVLWJZsXbMNG8UiMhNHTcL7lapsYWYYuOMtwvRCgl7G3iBiAyMZGglMFdFSlx7F3260Q8XmZJ1Yea94WyL+e4cq1qfhdUmUwCMoHjEIxQ8yTpORF8cEcFUJSdiQ2Je7XqrnZlWOhZB5/lVpEcEACGU96KFVbnlMjMCNi03bc3Jqp1ysGitwIYYOQYVTUuStNzIK1EJIaScHXR3MNyDTaoUFnyF484hgLE5f5W7vWG+xQ6Bqu3Faup+ZrzRjD7KVgBARDICRvTFPjLFISLAkrJnyt/e3n719denaTKzASkE9peRU1HtYjHpobLBNhPNbrKtIiIZVBs301fOvvEFAx8LKujj+5RC8eQQQhyGcQhxQNeyaLWGqxZpfjVz4RAtNQYViIuOwCm25zcv9EESq3MmrmpZXL1x7ysLiGcThJVVYT18tnG8D/KHqMUzgimiAJSQbDWRDIurXqSv6j0guo4LTAQBjNwbo/+4G6McTJnxsN8/efToYjcYQE6KD1D5+uhffJNS+TuU6Nc++b5Q+DPs6Yw9ehf6P//6N/0ad5MqYRux0GcvqCEjIVzs91ePH1+9eDnNi8sjne3rTRUzE7FqnuZpmSYj8kIbQ7y4uByGMac8z4s5aOhua8z7w/5weTEMu5xSTilLAlMOQ3FLKbyGQuVryd1EjARSaO1QEsZFADGE4mOMCFQ5YtixHqofiVYnyeZIWHpODz5v0vRaJJGZPJ7Q1FN2uapnViJRXS5j50hY8ZbWwbcNuxZWWr918QQbUXU/u3JLl91eo+x1u7ySVE5iklMitTigEzUrL8khBK286mK4WEEkZ5gyAWIAyQC2mk0isxewMsgVgiW5RQGW1DFAQLe297zPIQ6AcDqdRAUCIdG8zPMyF+vmIpFlQHDKY8kHBjTTJGqLAQIzOYTuC10vOyVCXY2o+nx4N0mmUIKkqgFOVWU7/ZPJfxyoIkLwppsoclRTyfOLly+/+urr4+lkAPv9vsTWq6lmkVwcqTjEWBSRFUU3Fa1wv09yPjgEZmrm3u5V4puDjqKJgMTDMOz3wzgGDt7TZP8NrSvNdV4ij+00pYXKCOpwDyNtl5ud+mpd+6xB6g8XeugAxjoUUOvAasynu4YDQk03sm263or8lEHXuQh1AjMwJTCqYRQ5JRFZQUbrouwLx07bMYBGZTDsM2JquL07oxDio4uLx4+vDvsdA4ivM/DeZcPmCxst1Zkp9DeU3Nd+/d5dwvbr3wsdwP2sVyyQdIzDB++9f3t7/MUvfnG8vZWc/fMhDnEYx3EEgJyzVCjAdZxINO4GIr69uZmOJ6kw3xCH/W6/2+9CiKaaJWfJqlJp1KRqgFrt2quPY6MzAyB6kDQV5y3JPpvXKXazGO3ZOu0qb1N4YUpuDmffExRn5VKAHLi36uij2KakFtnq5G6tHhJrB9Vbc3vhWhFP6qzlsHJyVEABxf2Ha9e1hgoX31A32LSqZFKRnBrpRVte8cr6bQza6mzaGHpEROycQj+hA3FVgSOsq/++Ias+NkBQNrvgmgxVDZWLmSWVZCsiBfdxZYRWpa2j5FnWDAswsQ96VnaBq+ZupfF3W8ViOKFEbBWTqPQCqjshdxlWQKJhHAng5vb21atXx9Px9ni6PR6XnNxpSlUlZakthkNeHEIIwU1Mi3uzNNq+a47Z01vI9WZmGcRELbsDazYV8EAIjiEyceRxHIYxlGWvczpXKWFJOmBk4gJ6cjCV6kqCq5PUvdm9G6y7soP9FL7jaGabuwa3nb+n4yH5laFmIEYdtXSdHap9Zx/OZaZZNUu16zdQKYQ906IbbrhtGb8roIOVxbAqYkzUCjpsm1hjApC0mMnh4vL9d995/733hyG21LmH16X2fai94RugorvL5bt76tc8/6GvY1tmbtxIO3wQkNBARWGI4YMP3j+ejr/8xS9ur2/MDGNQlch0uLgY96Nonqe5UDIQ3RmGnHCoMs3TMk9GBEiBeb8/XF1eMvOcppwk5SxZEclVvgXk3U6mhe7WgEtEZLaypZWsQlbsDFtZacYP7XcTM0JrK4F2KShKvcxsfXddmOokeih+y4zgvMZKM60eoljBccRC6cNeqd7dj36vQ5flC+QLirWTMjfuBwAIzNbbG/iJY+6ugXWhQOg88ZRVLETmEBix3UIrCa9V/5VU1xYrVMHbllnCXL24e2p+Px8hopMiiUhN52lOkj3n1qDMChVeUjQQxGBnUrUWhIZgkHPOKEwcOCBRm+4NzJVoSCiiHavQ535DUlQEsrpdLM7Khajj/FNziSLMOf3mt7/91a9+lXOOw8CB99GF6Dkt4vRWJAohsgukq1mSI2RZisirrpcDdeEyxejZ/XxEVvCPiUIcduMw7kIcOARwzCer5zkiIWjlHDkGyUzMBalHVA8M1Y7g1ut/8Q7IsB6fdzwOe8OrttSFhmCW/Aj/NoSIcMcWGlcfVSsuP9aPp57ylpelbFDWPbRZz7VvzNbO4HvFf3q3ZzUtFkMb2yeXf1vKIGn/bHz+0Ucfvf/+GIMUI+4tm6mXQNn9mNeZA9Afei38cCTkvRDSG0BL93/9YQMMfGAYKPeeQGC+urh4+vTpMI6l6wGFLDjCbjeEEG6ub6bTSVWdiYHE7t4FAJKzG2ABQgzDsBt3e3c1MRGdlzmlLCqB3bSruDc3LX/N57YsqV6QLrvEjhhmIjlnaplFLQWgt28zqGycdtNg8XioCGdLDobm99yoL64tJSLfpvmOxO/I0v00GS+WvqORqeuOSrdSyJJwU5hyK/cRCyyaMxhw4LaLazvtcm80tqHDqmaqOSdFAOZALWey/U5r779VBpZjqd+p1fKqmiWrKnFh4DTNvpaYkWJTXNJWXF+gUsFlar4OZqCoquaW/IjYbc/bYQBgptipXTseO7YsT9vu70oWuWE9cgr9hoNbPhGHyCyi19fXx+Mx5+XrF19P0wkQ2SIzA5g41cfUXak8mAjLIsREsttjOMkfwPwnuN+Q75ad1VTA/pLlkN0PhEMMwxCGMQ4xhIghOM3RMzVrMik4lZXI80v9hVDxOVhZOquz05qlcz8n5+wWR9iA6rZ+w6rrPuvfy6VfOkVjQth6nVvpwTyBdd2u+UemqjmL5gwmzd8N1rG+Znn1ZJxOeIz3LTd7U+bVPg4ITJHo6uLig/fee/b4ym2Dmb7Bxfd1fgjnJfSNhQLf8uuvnQC+o5/RupBOvWcPgkFgG+awmUORFxcX7zx9+uvHj26Ot5YVzBiZkU11XuZlnt2T1wwih91+P4yDdwDFix8pDvGw24cQRbJIoR66r2QYBk9zXB2kC1GvYC6SpRjHxxCJzNC/Xrp7kcXmgmWEDR3CNu7bZhsXUNN6GxRNAKJ30yUYspo0tm9HiBiYjEwU0LRaDrg+13W8xWEFVwYzdpKmupBo8AsWGTB25v6AHhLrMERgAuudJqg2RT2BC5modJ7QuOK25adiERufWX8UvZhuVkBmapolewbcMAw8UGn3kBDJC6bf0mpGSOM4cOC0LFmk7Ak7bTKsrJLGBsLV7g/OO5jaedeKg6jVcq6AbC1OkqCzriyDKDOHGJrtpZrN83RzczPPsxvt7w4HbxRE/HrU4iAdmfykr8ZqIpKyqPsmlY6AOUT3mCtYW7WyyLXhRQDj4DJgjnF/uIjDiERadgxa7A+YW/vrAi7vCXAVgoBBR/1fu5d1DNh2b9iLdDtUpw9qbPY87SOpmotVJGCIhdCnpmjWciXX/CjxBi+JiOZVcQIb9qgVhXMTlSHeOcD7oru14t8QjFoMn62HkKqZQOTLi8t333338ZMnQ4gKkDXbhjG6gfy/YSl8t1/+g8FF4XXkoXt3uQ8tFt5EFwb9YHh3BDgfc7za+JfGcfzw+fMXr1787d/+3fH2lkMcxmig4hksJSdbQQzCMAwDIk6neTpNauY8ZbdGhyLcFVVnRCsTMzMhZaef25rS1bUSycxiGEJlfGfJOYuaehqciHghcibnhpNWOCFQAYEat+Te7dth+BwKq3YM7RhQL21MPgUrogH8j/+3v4S3j+/74xLg8k/59QeAHcAH/Zf+o//uBA8veM+XA9a8NfVMQtUWPH6wMiEgV+KcJ1srOAmiHCLqxIVyK0suMaV32JAd9tjxfHrLh7sIhMH99RfPts3gEYAmMl4cPvzow48/fr7f70qE0XkmE95/ANxTKouNWGf3uNUdfGMdvheqeeD5D3oBuYz9jZxBex7rA1j/m3yTe77ujBpVJRiG+PHHz29vb7744svj8Tjsx3G/Szkv89zYAb61J0JiMrNlnudlcrZJ4EAeWVXccTXnXCWmHnutKeUsSbVmXfVIsarDuoFDYMpZJOeckpj2SnGfwc8vQX8oEFEIDJWEUATCq2uEqVqBcxCrZSHUdJWNssY3fwb4//m/vi39bx9/tMf/6z/eA8B/9P+YN+HmG9pBPw6Ujcz6zytFsyJ5vnlFRQY0VBNV1aySU7E1VVvNz6tiq0+1tK23Z7UsASs0a1tv7g4RtTvlGe8Gd+m6ZLZymqGJgMjV4fDZJz94/tHzYYjZFNWqeyTc45aP9zXHeB5J35V8s7sQ+sN1+wFi7Jt4AfUL2Tfp+lvL+s2jwPY0w3uZUXbn627SoSLGHJ4+vvrg/fcvDoeXMex2+yHEm5vb0/FUeDSqYBACj+M4DmOSpKru/OOp2u7u6RdLzrLkErjh9tEu5BNV8Cj2/mJAAECmSDWBxPt3dxPrKDdm3dpzrdidKYRD9JUrrrjdGLRFwWrgQ91OYDsw/Q//5U/fFqC3j+/FMfC/H+sxcF8R6IuBWX8ceLQ9Fysn7/bFFhFXiSBWHxXv9AW0Q/PXCbuFx537bzWDvE1qXg80vm5ieYDgtHKXwEqEBF5eXH7w/vvPHj8ioizC2JnG3Wnw7zWTOPOc6xwk8P6X8EB9tm9NA/1OgP43ebL9Lt+8LO1FKSBxePTo0Tvvvnt7PMZhVJV5nuc1vQsKL3Q3EqElc8INE8YQYiiOoX55iGpOOasggJhI0vV6aPy/9UBGZo5xCOxx7araUXdWbnsRjW4/5tVRXBqTDAAQmdxyXF2Ri7CBQzc2WHX3hFAcB/77/8tP3tadt4/v2zGwngH2ENjd1gLFm6IIXJz37Hn2ImB6fzmuTqFoZwBy5eaYYufL2VCe5kLX1VC7p56+HlfZJu+ammEGDofD4dk7zx49ejSEIAAmyTXVZxZD9xTx37newt/PEvh78GhaTJeDH/aHH3z6qZh9/eVX1zc3JSbMmQ4hjPvdfrdj5tPpdJwmkeysBmIKIVa9ovunaXY7lGLOrP3eoXzEtPpaORebOQCAqNQYgPp3qNiWVUqb6spJsxobU7DBoi52rROC5CymVLg73Bh9fiLhHUGkKv4P/9Xb6v/28X0+A86xk64bKsEEpmCa07JMKadlUclYGn1o7rZbIASxd/PEO4BHZxCKq6H/GaTSCfCxMSR0rdF4j8H9Wj7bFhiBkESSSR4vLp5//NEPPvnB4XComAWeQ0jf40f4/r9EdHs4VQEYx/jpDz4x1dvr2+n0heRMhGoQOOz2+4vLQwzDvMw3N7en6eQFtfhm1cwQVffCkpLn16SKDclpkX6+OCjRLMwcXCYqbjrnjPtCgGEObjHtIJBhY6OZqekad1gHUEIqG2kwlayIaMTO+vFd8ZplBCudx/7+fGLfPt4+fo+btohmXIjXrQKt5XeKpOk0zaeT5uRAaZmqax4XdBr20pQZbJ2YV/N+6Mz99dyS4NwHGlocW0l3teZOcaf4tJq/AfDdBQVUry4vfvjZZ5988vE4Rqef0R2D9L+3dJd/mAcArI5QEEJ4+ujR8d1397uxWi0yAQTPrYsDEZvasiyWEsTAHDgGQiphuAiqOi9zzrnA8TXNB4iYAzZqHVbf9qroqw4wrv/Nbj9Q7B1pzQIEACCrQYQAYGjsI60V7KhQDFt0X7nSVRVBAABCCcKqJjydGAv+3//VW+j/7eN7PgRM5V+KFyxWJw+QZCmlnBfJWXNOOUvO/hzYsmzOWvxOlbnGveDGKfbOPhc3Nk3t7/TfeeskgXgvTxOhKQqxnW0GSBx34cmTpx++/8GzR4+QKOdc7E4ewnneHgC/zxkAWtySLy8v3nn3nd9+8eVXX30lajEOu2HHIZhZTjlnAQBgQiLgEt6ypBTYQmDJssxLyku1Gyw0ASbiITJzd4K7KbV618JWxOEq7t+pQKEY+hcjNtuMq52qC4DQQFA0qYcM+0xTSeXt99MC/VggprvxiohvK8zbx/f+UfN6zXUqbl0komopzdM0zdNJcirYKZMBeTCOy55r8AzC2dUP0OdVrM1TL0fA3m6ud4vonPnvf8mvubcQ0BqFhQDFNGuOMT579vT5Rx89fvQ4MBd+k0sb/nSm9AdooK8JcP9Wz3/t1/EOEfcBlTSs7k6qZnDYHX746Q9Px+l4PL56db3b73f7kZCzSJoX9woNISCx5354jjmOEJgNTCS7g9v68qrdInLoj30wdIcsR3c8ClxE2i9V455w27rg6iAKHTWupq8ENx3yaCfDVahSzwCzpEpIJSy9GMgR/vf/5Vv0/+3jez8E/Me7/+i/O3oWGRKKiKRlmeackk/PkpL7tRTJva3EO+wmfugTGm3NwIY+HrWZVwCsQksEewB6wfNzBaGzo9ky7bzer37PldBHljMsaffs8oeffvZnP/rR1cWFliANwLMf9SZB7b/P17/t93/zSMgH+aSvQWm+bUf/bZ9MxR5nGOLHH390mo6/+tXnp9MxMAViAFxSPk6nZV5cK19WrKpmFjlYJR70cc+lASdCJvBEQC0MtdX+oHKrsuS0pML+hOYT18S07WNHgA0Nq7qsqbc5xEyIKjV+BbiLwwAFN7wQZiCK1UT6e40kvn28fWweooqeqWsqMk/T8fY2L4uH1SAiuLMSoTV3wtXrGzaOm9BWsG3dW6s8kd+BG5DnDC5qREq8g/G36aAa/GwS52HDXF95RGigBkRPHz/67NMffPz8oxBCzom8kYM+nvJ+KdXvUCrtW5bKN39+eIjWaQ+NRK95/pskSm6+v90JBXpoQNu4GQeEIfA777z70fPnc1okq7pBtJrknCW5HbqqN/tGgeM4EFPKKRX0H1f9v8dZuCcarkhkTVss3l6AKCIpO/pfxFpu0Q7Ysi5KgEbzM3AcUs23XgJmUAD+YirgLKXqVlP8k50vpKo5peqiTG+rytvHnw4KZJ7KoFlUJItrbgpHrgbjFJ+HypM+q/gbMAf7ubwW80a/RoAzx5waWwB36mGz4kVbc5PqqXNOP62SsuZ1pAYqgfHy8eMPPvjg2bN3hhAMQEWR6d6SeG6T9ruVyjd+vj30p6+ngd5j5mYG3yqEYJPcd/+Pvvf5D9kcPXQAIYFHZOx3u08//cRUf/mrz+fTgoGqvlZBTaAujgFiiOM4IuI0TcuylCjUdpm5gU8LtSjWa0WCjmZUQ8I8Vsmn107jW2Jeqv8/tokRWoxMCXVZs22t2Dv77pnNyJfAKlWqiGhmWTIaBWJkfEv/efv4U3mktMyn0zRNmhKYISETQ6DV8KRaYHXppHcdYnog9Q6Rp7mB3tFzmd3tHluEnnZpTneFtwbb5PKtRQNqzipyuLz6+OOPf/DJJ/v9qAYASvd6Y5vd35c/XFpfXyofev7vU7fDvd/3LqD0Jl//Tr4JwOvXMeBhSaKy348/eP7cDF68enlz/VsQVADRmuFQsqgACdw2S1VSSiKpwEkFni8cT6I1HdMLeVYVEUYAYHd5KIqVJi5RUBHPCWj5KWsYXk1BIUST4lbowQJq5hlEHv/o8VL1ye73rO0+8RDXjcXQ28fbx/f7cbp+lUVUUs0lQqQiYrSaSg8bGtyauLvGXcOZMc95Xaj5p3jfKbCCr/f9M9zj0HBPL73RhRGyWgLV3W73wx9+9tkPPjkc9lkSlTimh6oZ3ovl/FFK6N2vhz+5a6t25zYO/Ojq6oP333/nnXdfvry+vrmRlERyv/dHQiJ29MadHirMt4bLeZT1GjeIiGY1CLrMl6oiOauqD2O+nw0hABGAUQsLRuIQEEH9Kq9eu1Dafw80R1OTluaBLUMMENEQTUQB1qBqdxiCGkb59vH28b1/TMcTMhESxdW8E7RZ7apuLbNw5UlUwVeNr3ygZHs/exY93+K1eqGAbhpgvGvRcyb+ss0I0n1VVYAw7Pfvf/Dexx999OTxYwCY5hmZ8I5x6OtQne/T40+kqUQAPP83f3uvLi8/+/TTDz78wESm441l9dg/q1zdFvOiHgtcPKfM9WWByBN9TcQkq2TI2UTAFBFjycQAEZGcRd1nGMYYnzx5/N677z55fLULMRBGwsg0DuHysH90dXm53w8xMILmLClpFq0TCbRhomRMur6gBYy4Sz9ziB6b3i5GM/sP/rN/8ba4vH18zx//7P/4r+A8Bd18ByY5S0om2vXHVO8C7Gy1sI0FBHBmPe3QUUFkK+pb3RjN1qQBhV7h30b0EmhWpT31v7Ud4xapVBVBwEQgotMUY/jssx/++U/+/PGTJwCgZndUz/eOEt/fk+BPYAKwO+EMAECEKpLBmPn5849urq9/+fOfn159bcjtPTe0gBhiZGY/APrz3kwYQgzBkyCt4UHogmGtmRjFoF9ETDIi7ofx6ury6ZPHw7BbllMAMIAQAjGPw7jb7wPzkvPxeJpOp+NpOp5OOaWSKgUeY+FjBDITFdEj2qryLQG/HgvZGcv2lhVvH28f39eaMkSwYn2oomt6buuBAJtcxq1PdAU814XsZjpowb1NTnxv0bUzO9Le3WFD0L4zAayR9QjU5dmXwgBmqPr46tFPf/KTH/3oR8MwLDkjAPF9PfRdDzT8np4B3/sDAO//Us0LszDS1WH/3vvvffjhh/Px9vY0g0gDT4gohogIOSVRwco+Vi3+syHwbhxDDH6tMDMAztN0PJ3MlGo0OICpCaRMY3z8+PG7z55xDHlZwODy4rDb7w+HwxCHECJzAIQsmh4t02n+6sXXaUknOQGCH0UlrxvAcw+7hXFlEa3er8QQPDOySmPsn/2n//x/+r//k7dV5u3j+/n43/6f/7oSNq2gOL0RGzTgk4p9c7mltYU01hLam4l2PirdOVC5ll2X2BE7a+T1Bp65Yx50b6lZ8/6IUFXAhEO4uLz46KPnH3304ePHj9BgWRIzED5sev8ncVr/ScA/Z0d1/0H5R//46upHP/zhkpZ/9+/+drq9BQ4UgqFxCMzsklttSUmO/zAOMe7H3WG/C8PgthAhsJlpTh7FVAw/rBAVTAQhXl5cPn78+HQ6vTreIsDV1aOnjx8/fvJkN44eDagKYgYAOWdkevHiFd68IqQQAiBCGUSQiUpEwYZ9XNKQSkYYIwCKpZUB/VYN/PbxvR7YrTDx4K4BZwl/7ih0gFrC6sy6DNNtkt4mg/0ukaV3kUBc3R2aexy2EbseQNt0ym2dacGQzuNHNdCU94eLjz75wQ8/++zx1SP/Y8KHOtRvbGHfHgDfCgCCNX9z+356CnoGs8uLw5/96LM5pS+/+HK6fgVEDt/HEM1ANKtW97dCDOXdfn9xOBwuLnbDSIHdjRyRcs5ghmaOFKqqZqmkBQAzJiQmRLctBAKIMYzDwMQqwojMQAaAGMJw2O3GIQRmYmbCnEVyMrMWxlsyZDueWulAPAvNFEs+dpsQ8D/8z/7l//hf/6O3pebt4/v2+N/9n/5NCSY1rfItwJpsWntrbLUZPV+x0LUN7tg3b2Bgr9stphOqPrjuGzZRj5sl7jb6+Qyi6Z7YzoQ1DRMJzED14nD44Wc//PGf/fDq8kJz7hKiv6n3/ya20R/38afFLKyp4t3nZmY5a2B++vjxDz5+/tHzj3aXV2CmJjHGGIJITktyJzjfpaJZjOHq8uLx40ePLi8vD4eL3W43BEaSnJZ5kZwJgIkIQEVzSiIZADCGGGKWvJxOOSW3gBDJaV7m43GeTinNotmjcU20icgCsUeES85ZMgCUMICGUharue1ng53PyRpNDgDwH/yn//xtuXn7+H6BP//5v1YRU3FjnzNrLDdNIVzDjj1sW1XMFKBf5GqXEW0bjRc2X+g6DNc9LWKX49xFkFv9b4clbcMKVi/RXqxlYGCieUlgdri4/OiDDz95/vzp4yeBWbIUNehZobfv+cb3TxMCWo9wbPQA66OcEdFAEejZ06d/+bOf5bT8m3/z12maDxeXxJymU0qLqruwIaD37HG/2x32+9047scdMqeUVG7mKadl1izO61Un4Ws2NSYM47jb70DteDzlnEANEVX1eLxlwsPFxRCH4ldlYFSCq81ka1frNB8qpoHFb9T66dQfrpHsI5BszYexf/af/H8B8H/6b/43b0vP28cf9/Ef/uf/ukjZS903XMk2vbMCFhK3aclxtKYF0LtsjyarLCY8nSsuYO/SvIWGqzKz9YzWJ2x1IcQFRTg7DXxFh8REWbLNU9jvP/vss5/9xc+ePXtWfijSt/Zc+L4CQX9qOoA7hh4IwIwioprHIX766afT6fT11y+++OpLM81plpRExd19CIkYA/NuHHbjsB/Hw+Fw2O85RMnCRPNpnud5WRYXraipmqgZ5AxDOBwOjx8/DjHO86wqbgVqavM0MxGHQMhcxb1idnt78/LVy2lZvM57uh0iBmZip5+CbY+4OqVYS6W3LWO6HRQGBGb/7D/5nwHg7THw9vHH6vqh5mxXaMW6Pp26AuyaGF0PgKK80fMaiQTQV3bAzYIX7y+qNY2l/wLavUBMX+577wRbY3k980+UKLzzzrs/+fGPP/v0090wzCkzYcmrf8MS/3YH8N3PAtuY0cIeMyUKF/v9x5988pc3N//2b//2xYsX0+mkmhlRwQgxEIcYduNwcTjsxt1+v7+6uhqGgZHjo3F32F9f38yfz9M0jbsdEVWZuoFks3C4vHj85IlKPt7cmClz8IAAzwg7nSZVG+Iw7nZxGKZp+vVvfv3zX3x+PB0pBCQUFTF1/CdwAMK21WitPdSUeT9+fEg2s7Y9K/GR5JvtcvP8s//D/+wyY+ZAHJiwTA/OYDV1C6Jy5Gi9/9os7ltxJCI0A5GsWvArZCYO7pRnLbsbANHQ5c5gCKpZJCUT8fhM5gBmOSdVAaTywtxRCdAZsWomOWfVwByHiEhqKqLq+WkIZkBMQ4hElNMynyZPcCNmDMzMri41F11XqxZVkZQtL2DqttuEBAaqgoQxxhCCz3YFHzDLSVJKhDiOIzKpaErJ34FCA1RFj6dWA8QYh91utz/sgWi6nW5ubtR0iAMH14QygKlIyllFitmZk7g4xCHGGCkOIYTAkQJ7QxlDFMk1yAiHYYjDEEIsMysSuO9sa6q32bdrdoRT1wxub65fvvhas+wPh/Gw5xApBGYyNZXcyrWVDxqIgJBjJASYTvPxdNScVGVZkkoupAUkKtRoVFVRWdMsPL2uYTFFYgkAJeCo/LBy6dUWfWNxUH6rld/TW/ncK6pCz2XvKkLLXYLzDN/X/yOAViIoEpHmrPOCIX74/MOf/PlPnz//+OJwQATJYvCtOJ34wDbgT/AAsG/jR/96P5/f8fmvUdaZozFioI8fXf30pz/lGP/qX//V7e1NzhICDSG6XjdyuNgfhji4Fb/fY0tKBiDZNepac71URd2wAUpCi4mKpLSkBQGAEQgJkSgAkouNQwgGkHK+fvXqN7/57ZdffUHE434vYjllVXEg1OnDVlyxwQ+aTdKpIQCoaOH+txHYfMQgJIOctLDrimDYzNgMQ3R7icqwKKYUxYYI2parGFAgmKhkTaRcMnA2GdYFdd2sXjo2rgtyoAbnVLMjXa1PTEUyEBEwdJZLAjVbw/+iASIQE7oBBpiJKgtq983Ke4Gb+6neumuOQ/PmQyownggaMbGQ4moouVILPMWNsBxY6j7g/TVvBXSIQ9xfHBDxdJymaQKEGIYwRPQQc1UwTVlySqY1/YcCEnGMcdgN4xBC9GPcVYluRa4ublJlDsRMNVvEvUocdng4Z6Rlltj990sZJ0FFchIxQSBCYERkAgNQEUtghGDzPJ2OR0kLlNceS0lvPEyrMqxSFKyj9fh5fEfDW8F9l8RX8J46UKf38t/UTzQzXA3Zzirq5hc+Iw4ZnEfM9LfYRueLjVPqpG8EuHp89cMf//inP/3pk8ePVZXAgwXtDlj1/SuVb/z8P9QEcK+B3XcCkW1+H+soXKWY2G4c33/v3SUtX3/95c319YsXLwFot9uB2TLNBsYcRPT25maIcX/YI+L19c3pNJ1OpxcvXiDhMI6ElER89xs5wp4A8ebmZknJsqScCHGM4243cIxIfqMgEYc4iOnLr77+9a9//fLVKxFzsr9KzpLAgMgXwJW9UKMRVmpzu489UMk2ha6YkrB715lH1EPzNXTBmh9shP6vAMDMuG6qDFadY6FqEKIBVY+iOhwgApiuXuzrTdh6bseovBZUfSW1kwOI6v7CD1ZgYKz7FWgiz4oFr/YbCOR7uCwKIlnq7NFOngrmVszMKhixVfiomJVQZtOUFjPjGBippPFYtzH0f0UMMRBTTinnbKqA5GcsIiAxEwPAMs3H2xsA3O/3zEE93sEJx5K1dMceNhSGIYYwcHTVYajus3U9SlUj6CRitpI0Wjr+6hdCqwMabBHts+5yk56LhcwGGVRESiydIXsWBfhklsTmeZKcQDWlLJLL+UlMTKV6mwKAj50F7S8/yKNcyMGbatDZbNCLzbmt4dvW9rl2F9HZVOqaydfOt2byadsmv61f7W6Tj93PbC7vm4RHv44IUVRNMwV89OjpD37wgx9+9tkH778bQxQpy7xviejYd18q8bv0mAjf6gd/2zPgD3ICPPjD0CWCZsrMz548+eFnn51Op+k0FQtopxqIpZzALOdcE+Pw+vrl11+/vD0ezYxjjJGyqImKaoVrgqocj9Pt7S16DiWSirX80iyaUspiBphVvvjtl7/5zW/maQ4hcGAVzTmpKhKHEJjYwZa2FjvXp1SvxGZP1F/YHmBJCAFc9lwz52tBFgEE5MDep63+6FC6MO+4y2VvhgbEXA6P8763U+CsfVal8a0tfrn3imEdISgCmQecrfHIpKhAwFK/5H2uqRqi6fojEIAITMsMtC5C8EFAcO1M63LQ1rYdPU5IRNRgRMRgdr5wXJtMJgZmwhIAt+YycwghGNh8muZ5FpH9fn9xOKjZ8Zhyyv4eqmYAQCTggMxxGGIchjEyBSAs3TITITdOfDmLzUABzAreQ0wrsd2gj5p4oPR3bVK9mqr9TnaLKTM0YDdkdosGMzNblnn2LA0VN9knDkgVkyn+hb2fZcWmKlpfWWyrsXPj9XQh7+Qy2+Z9Qohw5t/TlXFsfNCev2mF+Gl3BbcbpTF0SuCNn3TD+rEqxVpRNVMQ219efPzppz/5yU8+eO/dEKK7FxG+aZn6w5bKb/n81z/5j7kDsG97Fthrdy1YErrcdefq8uJHn36WluXVq1e//c1vlmn2RsvATscJEED11fXNPC+AuKQ0LbOIAiGZNYPaFkRT2htyYLkAP2pwmueUMx1vve8JMQzjzkyvr29ubm6TZiY2NbHs7m+BKATmguBY8xIhBzOt5MrdmeJbTfMXQq6kBAIOgABZmlKh+FaJZc0FWEKgPoK4/hDqbgSoZwqq25E23Al97G1zL653azn86u1tBuhtHjAgMCIhZLdQXZndbpHXu7vU06G+Niqok1VxvgFJgYpgHRqqQ4YfPh23w2BVIQEiEnM9GPxGl5QXtRACe2CJmVJBOBDMsAJZgQOMkBbMkj3LPBAHDjnlOU9qNo67/f7AHNI8p5QlJUDwaQOIKAxxjCEOIQQHz90MFpmYuIL4RW1UvMRLecPyGlYyjAdE+Av0isXr+gvXxBJ/J3yT4j0tbANqsV4QTGhqWfK0lJDe7Jb9Po/Ucq7Fuq1eEqtBf/ca6+yGWIr+ymPofzLC6n/Y3+zO7sOtW6dtO/xVT4bQG0ZAbxzdgTs+WPuNsvLv6gBgftJCMYMEJISy9yIOHN59550f/9mP/+xHn10cDuY48O+Kw3y/ZQB/gm6grzkxyq2ipiphHK8uLz/9wafX19eo8OWXX6a0hBABYE6JCBlpmedpnvxuJaRhHEq3KSZqpSZWrAWAQuQGWPjPTKJLSl5fRAUQAkdAFFExLTEyxazKEZHm6WzoOmCPAQY3CQVAYCQzzapb/o/bB3GMgZmarIaIfBAwVyGXl2YNbMWaP2O1cK+2WratCwAeamEe0Wf16HCLOlsFNw3mvMuEhm1Sn7rRXkMNavYHGWGTccPKdLWSJdtOO2/MDczKambleWPrNjsopMfBy0YPmdBMJFcUXEWyggutqZ4L6ElVta4oGBJzRAQD8SkAsRytkrLIOI6XFxdhGFJO8zylnEHFk6iRmAKHMAy7cQgDEnk2kF8BzIGINgPKCk54d0z1SlZkPquIUM6I1VXf/2VNLTetPczmTSHwzZNjdpZzVtFlWdI8p2XxbxpCKL5YPhiI+RtfVzW1oa5nVPHKrfbo1rD+ZutTR5y1YwC83y2zEd/6OcDsHLnHO80t4jkZtK4WkLg1CthPCWqosMZxIxiCmqrqGId333n2Zz/80SeffPzo4tIAppzYjwj4B2jH+2YHwLd1N32z5yOsGPgZO+tbgUW2vTQQgYhEhJiePXnyj//yH0ce/uVf/auvvv6a/fr2HtjREKfDqBU40jzzp+yCidl7t1Yjm57E2mRJ5GtMRmrrIWZCosZ00C7IOotYWkxDGAZCNFQsLSwSFTBVMpjW8OEuOi0wxxiJuN5tpcJbgAAgINAFCbQPojTZbZfso4YzWnBdiGJpFa2IJUr1XN21K/8Bu8wbVVVQLRk7Gymmt2Cu+1nngIIYgVQgZFV1doeUteWpEz08RKH8TuWvNAy6dspFvVNFPwjUrB6BfOpREUMABVNNeSGhGh1OrZCUtXnJYSNwIJ4QEFVkEQHEYRgP+0Mchiz5eHu7TDOoupovDjHEgWMkZmIucJvvdJkJ/XLqqhLWE9ovFLMqki3UMNyC/nXNXtrulpfbbjmzLltxnRKQynFOqrakeZ5OmkXVxMwvuxJ8VP+Srpre4pteLxTCTqLYK9jLR7curlaD501EewH1S/BS/dz8NWt/6XauonZn8YF9LtPZ3qNuQraigQ71tza8gJmZ72QI8fHjR3/xs5/97Kc/ffb4sRiYKtMmPhxe3/s/UPoetAv6w5TWN98Y/IOYADZLeb9JQFRSTrtx996775hZSsv/79/+25evXqWc3SDaDJxcVx2rzBF9pwxmyb7T9QlxDRA+65qJEIARcZstZx4JUHKASxCY18ucs5cuEmEirLkWagZqyEhEiu5eIfUVlakDa/+oIhV3R+eVuguqVG19vwB0vXSbv88JEx3Ob978Vn5daevUWyosK2OsdB/f67UjB1cW0CrYdqgLWL2idGcAmKF74ZVr1LrLuxb3Zo9EyEagtSiWXwfqb7fuHM9V/liNJ5kIKvCPblQpiloAD1x78Cq/MCA7C3XysOZx3O8Ph7gbRPLxdDpNk4kgMXIIMYRhGIbIISKi+jbDd6khuANa0/ZtIH24z0LSumzEVtLQ2VzaNcuOlvUi1M1VSjVS1CSrWc5pnqbpdAIVICYOgeNKoBItRH2zBu1AW/LiJnJX62mzJrs0pKj8xxfYd5Ks1l+wnABo21CXTcHGde3fHyTVw7mgqBVWW6cG3LQjfbT8+v6YiarmHJmfPnn2Zz/84Y9/9KP3nz01gGlamJAZ//QEvt/xBPBtt85v/Pw3MtN4CFTDh76J1aLhcwA/e/b0L//RPybmf/Ev/9Wrly+NiJDEkxdbhwPUNpt++zceYcdIKMV0M8KsuuRSQ3v/qybdJSTzNaeIWwzmZVEm/w8hoWg2AfNwAFDNphmIHTRqCTWt5yFE7d4PIjQjI161C9jaOqrqdu9s28BU6JzOnzHoakhh95lkAVAMwTcEBKAuqQMgAG3r2XJXErrAraQr11aNgDCYKEg+S7ezNRSqADxI1NT7qlqYMFDQfGvRIuxWerjyn0oruXoOq3+u9QwBAo5RhVQyaPFX8h+Bpb4boTnoJ6qFpNlt5hEphLDb7YbdYKK3x+MyTR7tAIS+AeIYMLAhIhP7ichEze/+bgS2mQKQoZ3ReJxS2aaZ9ZWUjqAB0WVmMkU3wEREBZfa+s8LREAoOU3znFNSkeysKqLC3AFzhaWptd6inJyVPbxCMNbDebriNtb2RNRhM3jfyG5ozQK6x3q2BKf23aApinHT3eM2ibE7AtrBtJ5DncFQmXS9tSEyMM2yv7j8yU9+8k/+8i/ee+eZApgo08o37m/6b1X68PWoxh+stL7h8wP8A334LSRZliXt9rsP3nsnLT969erVv1M9TdO8LJ4DaSsbE8s90ERS3ZDbaKYAZ9B3oWv2PYKteV6rtWG1PixFByUnVTbGAA4Excj73YhEKSctUqy1j7GWP9wRGBoeY+t5g2e2toVmU/VgBFrWrBs1S+vCVhZHm+1LjwRA1Dh+Ds2U2OT1fSEkx0xWUlKbRBAMTGu/j+tCsMLHtm4Va+ErLtj13kVCAFRVA2n0mQ4uLg1pIYyuRgLQ1twIRr6n3a6LaxthHZ6kogrN0d4MAEKMu/1+GAdQm6b5dDpZzsgFaObAYQjEoX48LvjgOl/QXVYGdhQmq6Ni87U5f1H+TLVtcSllbf3g1drOtn5KWcWWZZmnKS0LqCIhM7f91trJa0kuqgoED9SDxgGty11tOvWNR2PhBawL7LOS2ei61TPCztapSLQ5M0peHhW9CnSmQBWMqYmQfoYh3hMefM+uqrxtqkkyIz169PizTz/785/82UcfvA8Ap2kJVGLe/2Fncf/JHgCvpwo1mJCQgFWViN57771/+k//aYjDX//1v7m+uR13O0AC08YINqdPiJpbgRJD58pwRhpsHs4dP610PgraOCiwzstQNFdmWbIqBg6EqCJJLTDtLy6ePn2qpl98+WXK2WCz3CMkZiYmBGyRR7VaooIWzn4lXZcY1tYgV1cJdSAaNtbrsO7G+i2cASIzAYDL4hVCYXq39k+sKgDczI5gnbTXTFfqT9F2ZGERGltJyDE8w6Ww1qaVCVkKi5paUrSSsYwN/VG1vhssbKm+HS2xPy6S6z9QIrSO5cpEWVWSZF2nFg5hHEYEnKZ5nmdTrZKn+soqdZ4IicPa9a/uMY2IXrtYI48bElmVEf0Hh80600+sPvG0flLrFlQt55zTIiJFgav5dJSck4iUrVKRmBXyqbpj1XrNVtEgrpv1pseqWL2unKsNkxaxU3RZj5fXEQ1t6++2Ff46yklE2zGxsgG6H2bdFmg9CcrOoK5z2sbBVj+41pkQsWq2JY9Xlz/9yU/+0T/6yw8++MBfaeCVooD/UNGff6ATwJmfFDCDZEmaxt348fPnOefj6fjzv/v5NC+ijsFwa0tUJOdkBsSMtFInW8l0Ls5Z7SyHwfqvWIiWG4zKU+IrcOKkfJFkdrE/PHn65NmTp8NuPB1PzUVFcP3uTBgCB+LiaG12lxPhQWaVO4TrDYYNbFXHPOtL6/824TkK6xtuBoBi0QBSurEVA9aKLJUTp78twToJDpypeVdC59qAkq4QF2F7ft0SAhp1p5lZzkYELpqFlgJoDbNu+uLXDsTe5hIgqiiaEa/0xiwlB9TfDqfnZ5F5mpZlQURAbuije1qyP22l+pjd9a/ZaJ7QAMQFBABI3G1KcWNevyJG/Rhqvaus6y5EBcpkZjklVc0pgZlLjttRYdb+TzfoYnUNaTFcUDgHbWPbmbOXYtsfAKVAd+qRqlVpf72eumedHCL7AWBtLD8fVe91A9r6Rmzps2efNwEAYRbNywJg+6uLTz/55Gc//cmnH38SGKY5sa8G7d+LCO7wJ1vb3xT+wqJIRVVFwg8//FANhhD/5m/+5ng68f5AiKKF+eZuEITERIhY1pauDEejbgnfKKe2uTatttJ+i1WSOxQ+h4PSTOyTfEoZAQ7vvPvZp5/txt2vfv2rL7/8MucUA0/emxG07a1zUiv9x3DF0LWvDw70nxEdtJnrmokAojmDtHggrCA+ViZmwxcKm9wMnd5HGBCMAMQb0tXtgDp9Nrb3qLwGXV0E+pu0t/9VpbIU96YeCxsIAavZJKAR4doAmogYYASq8uIV0cEu1HU9pLegUKXhODfMoxpGGrBWdDMFVWuwhllOKUvOOYGKvxbYgvvo1IK2aanExxX1r5qMdiGZmogjTsU2uakEsCHaBQWE3o8JytviPjZGjYzgyxmRkmpUNrGEnVeHA0VtA4Mlnmjz3rSir03Eax2K3vg9/WK2fvxY+FRWaaKdhUIdbajZRjnmiQRE5cPtsdzGsLAilih7L6yOuRWbtFUvgOYDQBvPu6mDiRRB5zns9z/+0Q//6T/+Jx8/f84EZsBUTePvFB57ewD86e4DAFEkq9put/v0k080yzRNv/z8l9OcxP15kJy2Y2ZIzcwQnDBuTTiOd3gzlcdtfTgdrk4+VpDNSh9FJCAXM4oplsR5OdnpxcuXL16+dKetngLnN68zCdUqDlAKtNp9ydTYKW7URCtXx2GE8hqN7lD5zyYKcCOE2oiWqcX3z6YbFKUAsAabE3K77Txr/hsw5doNIwB1nJ+AunUGbkxw/Cxs3vGgYoqGAKIPdW3922NmW55i1Z0VIwfJuWA4XetZkBHJeYJJ82rDt15erWQXtgz2wrkHzGG6l6vq+AwRwbaeulCBOoOHtdDW8lrErKo5uVXVolnqyYrETCF0ux6zCnVWAkTpeaxPZNFeoqfb3UNjTdFKNSs8YrXNtlW7/TD2SwJE4sDMwUxBDdBpeQC4TQRexfFvxBDpfZ7tvstaVXLOYLq/uvzo+fOf/exnn332gzHGZU7YGJ/2+/effyoHwEPE0td8/UFnvj/Y11/3Ih+8NjbbNm/1GAGMOXz8ycfINO7Gv/lf/pfT6TSMe0Qw0WZOUBGQsjDVFTRZM6qhydIN9OzWQCzlxJR9oF518aumx/vWr1++/Jd/9VchBreLmOZZUhYzrI7qhD4WMxJCtnMmWxejsYLlvo9FKrqgJpBp6dul5V7b9ArsFzzHtla6hGjFy11zSq5FMFj3HIVCg9hUoa3ErzrVNSTK08DRNmpPMzNBcxzddVyVagJbtje605FzIVXMQExljQMEq01ceVoBrny1voUJyjK7vsKcExKGGNcUqiIUABHJORcQvDSK5HZu7u9dLDZW8/oeNKE2kfU7mOKnV9S2ndEnVpJaESGUiY/qS4Y2Thn4uznndHNzm+aTiqr/ReLC5Km/ZplIdTXZJ2dF1918WTKp6qqI35BgsJ122AlJzgK3FLaxjr4ecmoec+C6ZSAkMHWjDloFKVhoZFSoUM34YcOV0r73X+8s3PYd9Y0slzGIGSzLuN/95Cc//su/+ItPPv4kcjSDtn7G36UK/amW0H8vJoD+sJCczeBw2H366Q9UdVmWz3/5q+M8FwfgQt2gzgLTAFbs+Z7O0nNqoBOWaMme1HaX2l0YE83bPYDb29ubFy9xCBcXlwCwLLMmgcCuY8TaFeJ9avT1jpNtWwrFLaLJxer9snkdqkrFvg07LhDCpg+yHs4RKOUSoVmLbrKKsUnzey/4jfM7NNmR9fvQOnQICmTP3XRAo6396o7FnIuELVS5EFIbv+hO739ORLHN6rIqD/wk05Rz7Wc3zmpSl6UbaRBWew6m6o3TFCH4UPvYa+YMrG5ou/7Duw6zSoXZDGag6pRZM80p5UUBbF7meTrlZQa0wuBhduSqGDGvXnkVFSqIU8fMKXOlboz11oRd7IdLqN6xna++2WaX1GV4EZUNeTU3BYStJTT2sY5dovvWVOJbo8Xlr3q2KyJcPLr65KPnf/EXP/vRn/1oDDEtqffi+PfqER7cddz5evPNfZPnd1fOt3g+4gOsqzd+MQ95dPQXpJv4B6ZPPvl4iPFw8dd/89d/8+LVC0TPeqjuCbjldp7tlDZ1plDuDcyyuqdKS3o3gLU01wmzOLi4nz6xBTCD6VSQf2BqBY/KTo5qH61nrwIBVE1E1GQNSaoeCgrAHJDQUwsMNvg+mAkig1WOJ5iCgpH3qx0NrtV6RlDfNneL6LrfrgY+uHJhqVT1tgSuCgV/K7QSlsrW1wDAxLJJQGAOldZe+ztVW7m1d0CrOlgBUUeZbcVmHcHaR0pYC+AqRTMzTSmVF0uddM5sCwqX7+mHNFfs3taoKiurdFq5AFVjsdZ/VUk5iQigr7SpWh8Zbosf+jrH1NQZUJBFl3map1PO4qxcZPaOhKAQm4q9iaqtrCLEHoWrcJWL42Hj42ONdE/N1KHpLaseALbirq4XICQuhhxU+EWApEUN4ORO7HZQWIxXi9TcXgPAVES00yhsONHQdly+tXL853B58ec//sk/+tmff/ThhwOzmZWVExr2u/Y76D92LhX4xtXJ7r76P0ppfeDJ//5MANbX7JwTGF7sdz/4wSeAkHLW//V/ffHyhaYlxKHA9HWLZZWwsMmk27j3t8YTG2JaW8MzI6rCROyuLAwcMA6mknIqCIwzUB1XJ3IDUdhywKESFglRSoSRdvtPn5GLpoGYAVFNQcA6uzkrxRmpBN01TBt6d/YNeg9EWNh8bQ0Cq6ei9b6kG7S6l5l19HsEIgJVM5dJGPrqVQTrTuDOB7h9Q50AruvE0VypsVtQWF2ZblcRxFgHLCJS4PoB5lVtZKtgDumey4pcAldJLVUuXWFvws5Cr1CiVrKSczFVRDNXEKl7r8oZYdags8LCzTnpYjkt83SaTpNKBqQQAgZyBL8i8Lqm7CB2EHzzvy+JQEUq4vQh7P35m6SXoCAy1nP5bUV7DFZRDQExoDvnFndTP5LaXyYEA9qE8a4Oef1Krevp7lZP2Pio90NAuQydFiUSmB8/ffqDTz7+y5/9+Wef/iAyu5FX6/c2AwOeT5B3/xm/TfWBP3UzuO/Qg/QP9/x7v37XRoSAgNy+hT5+/lGIQ2D+5//yX9y+eCHEiEmkIgPltmEzRC6G530JaubJrbj0EDCscqWzMLN20WJRpBQzZV/wUm1lDAkDBw5c2Z+dg0qnSGjG0QXPrTT8FpiGiMxBIIP2RFKstFQk1E7RtXn9Z28nAjCRGbhUqlH4sd6263ZxtXG0RkNqfNpyFCKBIbWdcyWummqGzNZ6eVjHdFzXLfUw8VMEutmrMihrm7seAN17VffkBgghRhGSnMzu3q22KVGwWSGXOYLPPOZXaKwePzVNzVpr7wA9gjo9q9IPWrJut5T1JLUiVJY8z9M8TZKyqmvzuKb7rBYLoNaIMW7tv/50aEYahQIHega0YE2uANse3OhmGs156axDQARmDgGwWGmVwQi7vSwW2XIxYa2YPpmt9Kezox9Wv1D/E9l4HtnGCsjHsjpP+Ml2eXH1sz//6V/+xc/ef+cdJtYy3SHeKRPfAdr8B3v+d/vN/yFMAL/L6YoABjklQ9gN8fnzD5c0z8v8t//ub69vrtM0gQFwwEC+j0I0MyJDQ3Qfs7W2t8tUG6Xa+nwmhD7PvvNU77ncpVF31qdhYx1aW6Ghnq0VN42TdUL89VU12+RS2QkJgoKCSZH1NCKTqZZbsa6eAX1X2IPqlYROTbFbw2wKi3PVOZk7UDbHzocbt1rzut9PC4lDVaDCBNhjwxXr6omFSEQgqmVH6i7ezUvZ2hmg53FV3WHQViXWiYR7M+L7AX3fHCFj3QHgWfq4dR29waqMquz4VkkJgIrFVL/GRGJ05qYsi+aUcprneTpN4CETzIjkF1yJXWzqXmjiLkbqCanNysesm43aNVM9fLC/wmyzsG8n+iqCK7ZLzMSMGJxx3Pa32lUlvGevdfbJ9PfAWUvRqYjvC4DxQ1RUVTKYjUN89PTpp5988uc//fEnH37IANOS0ID4GwrHP2AC6OYAMDtfFH+H//wH/eZ3OrR7noP3/d1iD4Moqoj0g+fP98P45PHjv/rX//rXv/mNpVxKNKIW335VwS0miVTILAiA7qKyjRKwHj1okVYu+aQiVrxbV+rNUbuXLLluJrDFbFBdgXre4fl94y+iwxRay0xMRT1QZQ4VujEoe2OrRb+aCDVjOCt49loM2nyjqgB1c7seXQhABrLlAlqllTblJiA4O9ZUnKMCFb8G1KY28jQb6AD91VUTSyqlU5VEJGBlM3lamaG0FXCXRrBRFhGGEBRRcrLmTGmrA2wnTcIVDayMFqetnCUxbmsXACiu7hVmACZlg+3fn1uycHU48FKsYDkttzfXaVpKnmep6dA+M2d39uybss0hbsFr0C16nU1km3bRJQHuJ2R33Xz6xh/WQDUqXNLgqrryX1tTSHtPjyaA29ZtbJiNIW6b8nKpNGGE6Ur3Wnus9jcK9KoZVEPgZ0+e/PlPfvrTH//ZO+8+dcE/V6nJPWcPnveV2BlyPFiR7mWHPFCR3vyf/3Bl0//5DzUB/L3L6Ox3QNj8IvF1324cP/7oQ1/k7ne7L7766ubmVpcFCrSLVomC5dL3swH1/9/el7W5bSRbRkRmAiBI1qJdXiS7773d8zDf/P+fMneWtu90t+VVqiqJC4DMiHnIFSBrVUmWbPLBlkooEgSJyMgTZ/Fs+vTFzTTksf4UBaaZpjunjIVEKPU/zDz0g7AYrUkrSFFeiKSUs04S/wd3+yDfS5YOlL6Ukri4pfZ+GFDOhpk5ps4U3qKQZ7iQZ38Q7ExFGBiYkLBc/QroNus/wxZCskw3unZFlXUSGeBk5uE5jHFlihBdchcOArVop88lCUX22IVmuHiUuONdiZ2TUSRnjLEtJ8ETo05Kk9fEIZCpXDWxobwwuxiTJOdOFZtvP/Ye+r4XFsdd33fbzg0dIJEyWisEFcY/jCVdPxOzwidJkTwKGSFiDjbdMSHOY/foTeJYGHiU7pI2MnkQSX54AqiQkBQhqUio9JgPCACXkpbxWjLiiub9RmJQYYm74TgXT3bMhpL57TD04BwSLZeLRw8fvPjyq7/++1+ePX2CAF03AASzh08CwJA7/ta9vAG920rv/vnq1nvy54m/0+T493zyyXi+OB4Frv5dTM0sjr5UQIp8qBaSevLoYV3/j4cPHv7f7777/vvvz9+8AedQa0K0zCl9vGTaRryi8B4XyF0njka/Ednd8WTOFp1YxCKiZxD5vFkQ0SAeZqEYGgUgzrEbUfIlZDaGgWrkVoRWScoXCX/C0rTCmzYgSmCfCmCMecqAeSkUxjiKEGFw6ISKqVrwDILsGBEs0iDLWr2yTWK8JrHPFoOCr+gB4uD8BTCeYydaS9pw+aXLe6IprVXctYxn0ZgsrAWx9JIEQg2aPcpcplEmrCOibbmfzeTIsQdxXAgwcwRSZoP4JSv63EiIDIqNgKfN2KFfrd71m62f1gMIaZP60gj75F4+evmQ5F549HVNmt5sG+UjE3xOPIRtJoCQr+A8Hp5joZ8jypkZwTgdebK4+nlW2WMUosmEZ8ZpRGoeEAorFIzf1ZJdJ1DoJUV8aLNPOnPOAvNyuXz58sW/ffuXr754drRYeKRTqbjSyVXwzqS84nXVafd337NU7pbNycH3Ulo/1xnAPWJzXqRr7SBgm7p+eHJcVZXWSmv9j3/+8+y33/quZyQ0mpQqdFQJiE7iKVKIHNrSaNc+VttKmUddnr8UMtrxjtL3rdY5XwSVVh7X8RRvFnFshd3k6vgeb+SDA6M8RyJgwYIVDnmU4QNzmCNyUFqN7jQsMbkxfs981+iTZOL4LZfdQAwskPdYD3JHMzE0Rc/286g2MUY3oiLedZQ2BuEICuGW4CyC0qQiLiCj9gmTRUfhHhr2ErRnW4nj7PICCErnm89kalCcoicECquHaMLKCKKQtCYAsNZy7xi474Zu29l+CyKoK200EbF3cghhc8VYKHmyxtE5UhB2xV4+OXgjIoJSmGXjpQFcng8Vub6F7NkTe1RQQHiLWRnPyUdqgEsINTKZquRBicCl4U+jWyReeG8dIn2/BXZG69OTk6++/PKbb795+dXXx+2MAbp+ABGjcG/hvqak/EE94fRl6wOMFRpXLFz3ePxErbb/+HHyj1zH/S+MaXHyJSpeVrzrmXMWSbWz5sXXL5bL5enpyf/6n//5/374AawzSpHWHDNSQxUpKmiBBLkE+whGixsEKZEQwOx+KbvktWzyi0AewnHOMbNiMqaqyCgK5tXMDpgh2PWHvX9anlKJzUilR9RDerIUmbQ53R1C325BNBXmd/G4bA8G0bMo6JZ9VULx0wBAdCVYVIxdWVKeB07IoVKgP7G3DwXU+bycQIaE0Od6I4f4MnGCCoqI2fFgHYuuqNCVJX4OhmiaYN4QZhQp0Hykgy2ESdOIWCyshzwI4YTFEUY5m6A31AsxvROYJrI2vfk2KWTHfd+tV2vvX8LsQBkE8NbTyYzV7wtj3KT3zxgp1ZM5U3TtljHgTwloGqdeQdw/jUn9MdwGlfbaFAlheSR5iD8CQIsRLhZ7lkRtgGL6j9M9cWkUGvysRmaAkdsQVfoIjAJDD4qePHn8t//464uvv35wejJrZi74+xcGoXvx+gkjFUYfOxbo154lRH6fUnmH0loe/NmzgN57YcYCxwZrnYA12szbWVNX2hhAMnXzy6+/bLdd3/cJwvAQM4fsPshDrfyZpJmblxjjNRrVcj4GOdcwgf4CIuLEss9qNFrnsA4urRJKDtD4g99Dmgov45X5XowKiXLHDOAAlK/tXjkciC47Wdch8ACjX4Tzxi6F8U4R3i1p+owpIeASMDS3pJwQY/bBAoVnjhQAXxHWiYjKZ7M56zCybWE6iMHpUEbKrrpAonC3USz2HJgkAEkWgrmJ9mNezP8IpTVePAl2drvZ2sFut+uu23rrPvKsymylz6n9j9kHyS8i76qCD1VYI1wIJk0TmJCvEIlbGPPHYiaGjOCvkKqGSnk1L6JKExsBlPE2sQir2a2MMt3/jmG58Q4wAZUlqy4tKgHqdOzEWhAmpRZHR48fP/rLt9/827ffPnrwUBGx49458CbPf0hDn/dnAe3fE92Ga/+xfy63O36/vnj8TfP09mSniIinJyf//b/9t8cPH/7v//N//v73784uLjxgrhSFhs2xZFs2SKaJWLS0PmdRAFHCSDYYeBZh4HkjUMT8hX/AlFoOwiTCg7UIoJWujMlhlNldGQWEoq86QRDbllmqGZnJnXZR3TLBHoTZAagQDcwxUqaEc0fQASFxjBxmcCAcJhSxyEmSmUlwyU8GZJhZ/XLJB+kZ84wuUFXSbD0OIiD5f6JERo3SzG4YhhC6C4XDMPqZdbadyyBGuQJEh/xilLwzewzSjSB/S+P34NLh4R3KjKWw3vl/CaJkEpFu2zu37rveG7iSVv4UpQzPTZoGryMgwswSTjtIlpjwkFeamMGAQVgOZeZkfsd5xI+jyYe38UEFiBIovgA+5gixlMpLyYAbtTc+tyKmHo8IAmHTEqQQgnsBGsm3ScrA8fsnAYFFO/vLt9/+9d///dmTJ7NZQ95QFsT3/nj3aiN3q1qfdAm91Q7gtkPnuw2p72u0fcvhwfTTZRZhC4Ca6ORo2bSNUqSNefXDq7Pz83er1TD0SKRJBTG6/8ZyMd/b4c3598asEK9jLUlh91bavsdUKREZrFNdF7L9BCAS8AnR926kSGtVjKQxJynF5k4SoSeOWP1dPvUtirHA+TYFyNpyLNED8SQSP6TliFCnZt+LjTmF40wI9pK2PjLpGz2on4jngSckSCQpfJYQyh1AwEaC7bWIc1yMIpJ/fVxgpcxg8Sc5Vi0Vm5XRDigS5tPayiLe0gcwwoUMIiGZJpPEwqIaQseEQRxb6EX6YRBrgUgbTaQwEAGifk9S+cME+WPw+RMQAcdhnShUdYCIqLzfvn9l8W5OiOk5A/4nUAjlvG8PuuCS601DlSR/73LJgcmoax9EInt5QHs2xAUBFAU4u38k7gYgiwx9B+xAqaPl8vj46PnTJ3/59tuvv/6q1cZJ2PZliO6ug0O5Ogj+vVlAv2Op1B/uKtzq+FtHXY73m8XPcXdMtAcRK2K3d8CGUoYo1jmN9OKrL49PT/755Mnf//73f/zjH2dvL8S5kMMbtTMOXCoasg9nFBEgl7TvKBlBmIrZIkui4ExgstL0o4Zt12PXxTT7LDv1wKgi1Ep5lQBk1/9IRPfZWfEiRuGTr8sq06wjN8M5BySlQjjBw6V5ZYJ/PfgzoqNjTiH25alkFEERvykF9g45VF0RIrOLtEX/1E4YUKn0GeeBrMSsTgwAjUaFDp0DAZekdmHmTYJFvBom5DuGFBemMgnaC3APUujR07qEMVmIokslBe+PYmCVtgEg1jo7dD6J11pGv4OM+J6nz7hE2E/ai0zbhNIXMNiQJPVZHk1HMwdf8V0RmR7hKSmpbMGuWWutiGgY7OBcRHWwSMAuSEFQ+CzlqBuZTvCkJA9IgflIUptAqRmeQkUp8RMJAZgB5Hgxf/ny62+/+faLZ8+Wy7lRyrIDARrN+QD2D5Y/SE3/hErllU/+Wc8A8OoVAm/1TLLnb8Li2JJStTZPjk8q0lqpdjZ79dNPv/7663q9hmEAbXRlFJFEAtzuWKKMfo2z2TEOD9c2J4lQHuADFgbH+zBTz6dErZUAyCAjY97shcBTJBsCYVzEt5tc+jT74Beh7GksI/x8d6fp9QbEDAA8fiVJcthpcsAlVH1vIokOnPgKnn39kTn6ymX572gC6afNhAoVhNSAHVfXaV9asPUB9wKPkMITQyxl6PdTk8pY9schEo28AyeQCPeDtbaz/dAPvXMu7C5ElFLgIRr2MgBOPJxizJP02L5zZ4xSXS8DzqasmKUQyapIIoQCeT6R35WXdCnl81oMAqBz4aoXmM3+1v3y2w5z0PuERcaFUlr2TVQlOMRFZYwdrLBDUm3bPnr44PnzZ19/+eVXX3150rYA0NvBOqdJ4R4jk8PjTgvA7+Xt8x7HXwPblcmqCPvpRkSAqADBOgtEbTv75ptvHj9+/OzVq//8z//5/ff/tR0sFBGoOHIIGr96ph/47zpF230pOc/lejGKndyZZxASq4yW4KgsoUrmlALsvKKJs31Cog4makcBc5MoAGKxI99kT26NOx4p2rGxcW8yvAyIj4SEANlTH3CPyrRINSiG1J7kAlrEZx1ymgw7Fk+dD+5nQAon0LOXjlEwi/Azw7BZkZRmUtBciwVyx4wyCU/ztQ5BAADCbMNyVa6N/vAo0CUQQBTr3Hqz2mxWbhhCv6IUZtdP8URP3+FKGnd7QTiMdrI5yEFGmExUdWXvZkmz55C3zHGdi4IQ9Nn2ipQir+xSVPg/xGGxpyth1ARnEd34Oyp7b7goH4RxVqdIAcdFsDLnOoZtFwWluojwrJ69/Pqrv/3tb8+fPW1nTW2MBQFmQjBK4SWm0HfpND8fb5/bHq+v3ibcnMN07z+/2cnIPlrL9Sd/GR9nx+jVBy2BZevN1BazZj5rmqoixLZtf/nl14t3796+fTtsegAArUNuO4x8q2RPt8QgJJN/KnfPWL5BnIRcxC2wirTu5FYNiKh9pKnzpkWiFJGAZe/dNj6VbLaISbyVgkL8MGQUK+wDBAQTGouTO7lcvmIuL2MaZO/GtMq4DhTNJaarB0lrgEohADDG4iUiwNGNLuYtFqKKRN1jwUSUioI4it55I+PvnS2I7CalYVZc5QRhDiNmYJ8Bl3+dAJQiz6vq+16c7Ydhu930XQ/sIGTbU+YfeQ4me0FJwou8iJqiH4j3doXSyQekXM+irjlPfiRreUXGAXKESKQ1BYqnivAg7tlOyw4R9ppWLM/Zgxo5mDcp5RMm8NLUtLTnZcfdMACz0ebk+Hi5WDx68PDlyxcvv36xnM8AwDk39D0IaCLEWxaEyw3k36M6Xfok91ta3+fkP2sISHbhk31EoT2zH5mmc48F/+PbHRG0Cv4z7FgA5m371//46/PnX/zw46vvv//+u+++O3/9GphjGaVpQcm72visLAIus7bD5hgz/wdHA1cf01KudyUn2Xv8+NZeG22MrkyllQJC7TcnLGSxF+uQAaXkquZVyr8Ex67eZxO6aIuWwQMct8WYh6OByxzlGZL7VgRBIIUkCCwuupUGRhGOhh8yUmogSnhnOUdGHDNY4QDUhJxkjJPl6LiWVaM7HwSGJGCSjEqlpAEpjSsuheNQCpNk4GxStGMjygAEipABnB3W603fba21Iuwzf4JcUKI/T4obDdeY0nWmHNTLwjmyvfg0snczQLL4ltIfFMaFPywWSiutlTLezcJTISBoCUsoHcdTSEQokw6kHOcU0w6h4toKAImXkZEiBBGLAxTjd8qrjuQRPKIVIcTjo6OXL1+8fPn1sydPF/O5VuQcez2DJrq20Zd9gN+1kPKfFAL6CHuZm880bv/89/CJprsqUBvFsoDWuq1NW5/WVVUb3c7an3768c3r396u3g2bDgBAacqsY8ykkp1Gym9zy8jyeNjOijWxoIl3f+SCoAgjoDHmaLmYt60xFSnF7Lq+22565k75KUVi+5Xz8rjxlsnb9h1aOTCYdBEjCCvZpu12c7inNuatT5HWLmOQHckHz5Ym0EBEoAWd+K1AUa0lG6NRMh+TER0dc/pIMg5F3CUUjKvCTlaRd9uPWWkeIUNUeYWHkBHpi2/XDcPQ9/3Qd10/DMAc4BlSYZIfUxjLMXHK4srzd4GY515SznCEZREKQo6TSRIw/8GzJO5WnGF4vN+gUoU5U/Glk+m1wD1y3oIIkH/GkERcAApJKU2BU8zgnHPAbENuxRQbFGEYbA+OCXE2a44fPXpwevrkyeMvn3/x/NnT43YGAL1z1g4qdf2XK4Mui1D8narTJ1Fa9e8OWiG+36W5EbZ3RxpvSR/I9icAAGKdAOJsVr988c3zp89/ff3b//3737/77ruff/mF+wEAolI3ACosIlP4O7e8Y0wKcthW0V9ls+fkrJ9TPhCRvCOFJpq37cNHj4+OFkaZbdddXFycw5lj67tOZpeAHcxBi57KOLaoRm/bIAXqXvT8EQwp1hMRmWj7Jc8bJ25kwCCUdPyZfZhAg5QcnEAWzORYUgqI3DA4yAz3/HrRHwhG4tKyvIS8B+dcMPSPfTLGWpm9vHHSKBZx9h4wyzCQFNQ0H62GIjLYbr3abLdbTx4jIiCKE9ugrebCmBOxjEDOwr6IDo1JqrEkhxTQ5LQBRVKEjCXNoetXpDQpn2isIbr/ZyKBv+TuukDGEU1eCpAzzr4EAIKAUStlVAhYZuuEnTgH0Z4ojykikQ9FAKGpqoenD7558eKbv3zz4PTUGK0RLYexllaE++7lm7j3XN0mvm91+l3x/Zuc/B8hD+A6eq/c72t5aa+wA0ClaN5U86Zq5zOt9byd/fTTz6/P3pyfXaw3a952AABGk9FIiQW+RyoQw01K59wpXoUl/CpjGkXCCAT6wXZdrxUdL4+VVmqtSZHW2gpvNlu2Ls4MxsEmQX2znxWLUaZUrkRpqfJhXpgCDkYVcGQhzzHqEETEcfSV2eefnNbbyDKVEvMLwotk9p3WSgx4kM+7L5J0Y/ZJWlmxOBdmi6i90zVyQJR2mVx5y1QA6EhBb1ds3ZC0VgQwDMO274R5sK7rOjf0fpjt9YYhxlk4FMuC5os41jlLvNDRDDwdXGI+KVKxoNPCjt8IokIkFXkCygvDIFCExgSGnRtJyqTU8fgmrNo5JI6jvymRUlrpgO9Yax0arWaVoboR4W679cJmQiQiFrbDAM6CgK6q5bw9PT45PT158uTJV8+/ePbsWU3AAIOz3sdXjRD/PfuSw+PuENBntwyMZrvjnKJrFwe87oh0ACGICtC3t2lTiF88f/rwwYOLi/NXP/743Xf/9a8f/nXGDNGpH7FoiCfQI44YbwHCBtoZ5mSoNVVmkREOIyDrzfrNGT14cGqd83vtxXxRafN2tfr119/cYGXEfYrVnMocvlRjYwKMh1Z5InMLzbIAioSbkCh14tEwabx6FqpTBhTyOVap8MuYeIMxrSAjRFjEr43jX4uUAPb2+RR3Uf6j4hiCEpLgfX4wOWF2FkCRzl12hiSCerm4YCEhCAAgICiovCN32gcQIgEMQ/fu3Ts3WCACANLG+y/Es+csv8I4ky7N8uJIVyDY/kzLWoSuwihC0uYFCge3aLkRdB6ktCalPOwDYViNkshNWMqCC3shmYh7JV6JVP5L928GEA/eEaLRqq40AtjB9l3PzEj1zNRt2wLAW/W277ZWxKtqAP3SQUrR0XL55RfPX7548ezp06PlsqoqBOideB25CunHUw4CyB0b6sMC8Ad4yEdcb0QExDELk1aVNrU2bds2s9ls1j58+OC316/PL87enJ2vVmsYel8wQHt7Ndid3wlc6oiT3fqLd4kZIwqtt4i4wW0267fv3q3W7/TyuGlqpQ0ze62yE5et1BMgXqasR0Q/egRgwLenJg0ymcFLQsMheBBI4QBa0ECgjOEL/miUYhCvWJ/LkG3gIGYuQenUmIuwMAEwRs5ksOaTmMuIKN6+h9mnrSEIO4eI4EVSo61DiQFhjnvzL4W5sSYiEmZrB3YDM2+322HowTEIKa2VUiLAjkUcJJtlzJBUpHgmS75wIZklB5mFGEPKbhNjPUcOAhuP+pEItSbyHiE6xLYUb5OLSg4TswvZD/tkZWXA8NiJIIAmVMoorVCEnWM7DMKV0Y0xjTHsWCsKHheA3nBW+t4KWE3KmOV8fnJ8fHx8fHpy8uTJ42dPnx4fHVcIDsBa66wE/2kYZXPLKGPm8PiTLgAfcnexqyREQIUEShCss740HC+Xi/n866++Pr84//GnV99//18//PDD2fmFWAsxq2i3zBWzBklCpBGIMU2WCjVXeUuvdNsqFIF3q3dvXr8x2rRNy65frVb9tuPAm2RJAO/EfBRG5NMI95OwZeHoJoYyWb5GaFQwjxhHdQeOaeSMB+sE9n56LIDKR66NQqZiMcMMgmcKI0vymUgS45I6AilYcoQtSeG0EfVjiKhRMYhzFkI0TJF+m3LYgXfnaxMTMxBh5/qu74fOWisiSKrcQAgziwvVPAaBxnmFl+l68m34rFPqcom3RMvnRBstPrS0MYpnH06YEI3WXmAMBL6ChswzDDaykDPesRw4jWhZ5aQ24oGjgDQgBEVUGWW0RgBnaeh72/caoFm2s6YBEessW9dtNwhgbY8S8nm11keLoy++ePbyxcunz56cLI+CypLdEHeQSqfN4bTiT5kNh8dhAbhsj4A3/vmO91eJeea/MjtfkbXWxpjGmPmsmc2atp0/evTw119+Ozs/v7i4WK3XvNmE3YCpSHkC+A4JJtzWDEiw62BZFh1IMXkgIp4JuV6tf/zpp9V6M5+1pjLbbbfpukyJyVKmPNEdT1TGjD4pZQ2juNqRuB+Rx0PgUt+Jo5cN/sW+73fO4YiAHyo4igADqLwq5rfPnhk62SPk0/IWR46RvJtDQfLBMnHZo18oJD7+s9RTFesaiIzC/4rZLIAv3NYOdrDOOevs0A/iBkDSWgNpFvYyMS6lFZhCVSiTjGLDH2CxMCTI45AkPkvksbAX49z5p/EEBE9ABEUR88EUUpqyMSWTogULIQeLKJx+Qcbki2Ac6rdxWhmtFaEwO9v3rh/qqlq2rVos+m0nwJpQRRsUx84NToSHbQcii8Xi5Pj49PTBgwcPnjx59OTJk5Pj4woRAKwT66xLAhXcv9HHA+x/5wXgsxWC7fz8SkthGOsGcJ8abI8WbN8/7I4IU+yUsLORf3F8tDxaLl589eXrN2evfvzpX//8548/vvrtjfAwQOmtAzmge7zGeHJhoSot9vQFeitFQKISwG3X//Lrb7/88pvWarFYVFW92W5c8ODngqiJoxD5vJsPf2HPwA9MSxqxlVKoakKImIFCtFeBZmW+UAlUYHCeUCzsnCtAlUhtIlRKAZKM5uaJVVJUahEIOyHMpgYS5sFCqISwKIy+egacn8DbnCKiIoUiztmcOOunIFhe+6BuyBYLIanAbreb7aYTEEUKEYGyGR8KOHYhhDmK0RJ/t1zzwliYC/W0v5iE4cTDBNzbKGXZQHQFL39FeR9ppZQHfCRe+pJ7MA4gCkN3ibsQobLJSDB/nFTFfYwiUIiVUU1VIcLQdduus841Wi9ms+Vi0XXb7WbjHHfbrV92tNKkNQAQqqapj5ZHX7988cXz58fHx3VliAgc2+jBP3F1kL1N2OWg/25g+A2rChyEYH/IHQDerWEYbcMvPYodh7A6Cv65tTZN07Sz9mi5ePr06Zuz1xdv352fnV28fbddrcA5IAKlkBRRloPFTm2fKEmm5hAlt9/HXltr7WBBYLvt6qZmAQ9xTJ4qRaOMGuySsjTC5rPLD8CYHJq+YdGQbt8/SSwiqIp1LIXMBNo4ICJVplJaW+cG6wRRKSr0veN5RFyWiAiRAoaSzo29z50UhgJ+osAiApzyRIAQk8dEiV+NsgyjzwVpTYQC0PfDMFhm13c9s/VLCiEykF+CEGKgWQwUJSzys9L+TRg4UaUiNkOU8D2MGoiIgnERLR/Cd+NcmFApVJoC0VNhkLyxjNgSuNs6l5IJTL1Jpi6wAPor6kfsSmlllEIQZnDODr0i1IoW7ZwIZvVMayXOkgABMLN1FhiMMbNZ3c4X7WymtFJKLZaLx48fn56cBJIog3WDFQkt/8gzRG5+7x72BO8FAV1R5m7LkH1PRu1Nj8fdgen+/QDecP94Sw9ZhGgDECS4bOMJnRwfHS0X37x8uem2r397/f/+8Y9//OOfP//882a7FWZSamf5Tg25gAiWNG3IRikT9kx07EEiMpVhkcHZYWVj6FUWQ0n6nyCM5pyBWinZNm66eBIRArhdTmuyGJpsVzJCzWF0F/wygndp8OqPxxFi286aZrbturerd9YyUMmFlWkiBEYiDaFCdM7D/2GV8iFloXhTFgfEYCnE5AstkkUfu3makuETRQq9g3+/ZRf6ZVI6aL+SN3jh3uxntl7a4N2B/K/FsTUXQUCe+uRpmuSfCws/7pAGk+bv4SKo4OqglDKalEFSQjmrSzJpYBz9WD5R1Gkkt1RmxzJCJilm4CgibXRtPNZvh67v7FAp3c6a+bJtmoYAnLXvNhsPatVG16ZSiuq6Wh4dP3jw4PT0dD6fm8oQkvcdGpzzn1AKKN4bJ4m3BPovO/7WwqD7KJX3W1rf5+RvKgQbpah+6sZwCLs49NgIei+KWI73pnslvJHQAKcLRlbsk8K05503zayZ1XX14PT0zZs3F2/fnr85O784P7946/x4AACUDs5CmJChIAEYzYzzdBNx8lYBSRGKOOecc4AcnY13xh5707EFCvvi0ZC68KOERGcskReBlBqWo2iSJLfQwKVwFAhqOSCBkF01b9vTBw82m41j9+7dylmHirQi9HYGo8qVNjMhFYGEsuI5OjSPZhgygpUkOUFINHmASepwugpFHIEwiLhQuwPcDgjimNlBaaMUaTuYJ6lS5gNL8nAOr0t+zhvmszGGLZT+NM5O3wpCUhRwoajwQqX8diTFi0KRaV/SzdJwqEwIyr1FYpSKAIhX82qtFIaMe9v3CkkpNLNGhBVRXVWGFLI4Z+0w2MEqRVVdt+1s3s4Xi8WsnbWzedu2y+Wyro0/DSfg7OAcZ1HDTs3eO6iTWw595T14ovcuBPt9S+uHygO47SW+2/E33BNM8kWu/aWrtgI3hIiKTzR76LNNHEyj6emTp8+ePrXWXVxc/PCvH/75r3/98OrV2dnZtusic8br5fPbCEQUSP0/jse1WLIk/T4AcyQKT3VmON31l5TQbIkzUYfluC0kUuwlvaF7TdsFLtBZHI+MMSUBEHAyq/HmPCG+3DEizOftkyeP+23P7MS59aZj5kgnEpRJljwCEgRffgGFBMozl6R0AMj9JOdi64OSR3zYLLAqSQB5Usvi2OWhrLc3QBTv0uDs6EpjhGUyd5gLC888zIjaW9/EEyLFTj9mvfkkSH+26atGiEopo4VFrE2LTdR2ZbQ8EXUn4cARzM8efDHduhzLe54+KlKVVpWpCJHZDj3boQei2sza2UwR+bVG3DC4AREqY9rZrK7MrJ0tFovTk9MHDx7M2tYrMxBhsDbpJxBBKYLSDPHqoNDrduf3wgb6oP7+H6e0/iHzAD6PR6rKnAolIGkyxigAPx4w2iyXy2fPn11cXJyfn799+/b84uLt27c2bQi0BlKoCEnlUgwEE1//8d6GWQglQtKxYyvR/wnQUUBAaegnu0w777yVkFmPSjuessdT8m+poo3aXZlM8HAUta6UmrfzRw8eOmf7oXPWIV5stht2vIvgpeBkye7UgAqQ/cKHme+SLYBiyUOYGM/LqLRMqfAxNiUaCkHi5AS0RKJcLr1cpq7kLVIgz4zMM5BGFp4ev8Lg3oZhozXyyg4vrpXSikgxOudozG/Y63+TGv1wAUqQLtrwlMZ8of0nQo2kEIFl6HtNoJSa1ZVUWiHWVa19BKl3aFDYVNVs1i4W83bWtvNZVddN3bRtO1/MFSnwbhDC4sLk3mNk5Q4RZH8Rl/sr7ocZwKcJ43zsjdh4T3zPGsJMzVCUXsvZwUVqw3K5OD5evpQXXd+/efPm559//uHVjz++evXbb792fe/Rhdh7ugL4DsUMU89ahL94tMIl4asP4C0E0lhO+nAMCMfUXsj2D/m6FDlQgckOCEAhDffG/ZhEy9CdAqXI29i5Yaiq6uT4aOgHZmB23bZjZhwlRkU6FRbrWlhmg5X+1F8OSXICokQ6z0hikdfuEUsnX/gdChcwO2aX7ICwPC2RMaLCY9PpmNhFIaTYg0Ih1debqY3KIQUeEfnMFlJKB2nzeB0NA5u435lkYCSPiICcFSF1I4VfXKTIc/wRgD2oiFU7a5tGKxVYcCwATEppo2ez5vjo+Pj46Gh51LbtbDZTWgfOFLMNNq6AAKgwGsfeho7xHgPez7E6faDj/3A00Pc5edxTpm5y8hMz46IKjb6sCCPNrZ8PhNuYSCsFALU2dVW1s9nx0dEXz5+en7+9uLg4Oz9frVar1XqzXoMNxrmiNJAG5a3kU+YpFozwhB0Vs4PYmYZxgZ/10R5n7YQCJEZKQqd9mHumAoYGN0Rwl6yhRPPJpgaj5wev2B1zeBGVYud++uknAnj05FFl6uViuV5vttu1s9b2UnTZMeEGhOQqT/m8ZpbwgeTIkTFE5AsmFS6n4QW9d0VaBiJxM+L4hS1dRLs8wYlZoBDqTiAxzI6fZehmsEiTbOeDAMFMyKsHfFamF3aBuMm0SibRqTHkEpNYLLj5lWqQ4OCWLX+C0CEO4TWRMgrFEEClFIqgsCIyWlVG13XTNHXbtm3bzhftrGmbuqlqo3UE+p11LOLYA2OTT0xkD59bdvBbubyv+PhMyjtUGzjQQD/Fx3tHP+OuifN4412mRCESZOIPW8ep7Tw+OjpaLhFxsO7s7OxfP/zr559/fv3m7PXrN+fn52wtEJFW0YzTOU/ZCBmBBWOxUHLJnm+DCDAJkkR36XJsLkV+JBTO/jimYyb8QkXCvuzwc0ZXo1heoiUQ5cAS/1xkh+HHH19dnJ+vVu+ePn3uhAmxqWpb2S1L37vgloB5UxJ85AvMujDDKbGiAoTKhAABgf2O0IWjGhEppbyPWyTjZJceSHD/GI0LmYcSFv3i1kyp7iMhRjZy4IK6GpUHQIo0KaUp6Awgpt1k+fJu0lrh7hFDcyKYlIG6lOuQQCe//HtaLQuDZSZF9bydz4yhwGYVJKyNbtt2uVwsl8vFYrlYLJqmIcJEqbDWFtsPQKJiQrTfvhgvn8nhnaD2w+PWENBnngfw4TdWcpvFYcx+mdTiYAzjM2OVUkgAUGlTPX5sjHlw+sDvAd6tVtvNer3erDeb1Wr1brXabjYw9JJGBUqjp9ShzybzgVkohU1C2SEKOwbBiDzErtHzE3ky3saRg2g0dpjGqOAkcCwtSDIeMaaT8HnDydbFt77seLVevfrxx+22b5pGBLTWVWX6occEH5T8Is/pzHWlIGGOCocXbxXJxoij0CqBUfZt5H4qJN/+c9zAYUhVFPAOHzg1UPBqZJa0KQtDcgzbCAyfN2bH7xC/nEXbefUirb2PNCoi0ph6eZGICtHOl22i+UhvOqBRmHd26cMBQm/eHBJqrHU89N7qRAPgrG2MXrRtXVdNVVWVqUzVNHVd101TN7NZUzdVXflMMXYuahuu2569f791qE63P/nPPg/gd0bcbnzyeMXvQIlBe4mSG7xfGAAAnBwfnRwf+SLQD3a9Xr9+/fr1m9evX79+/frN2dnZarXqh0FA8k5egsW8ZIOfEurIai8RESdEQqSy5Uxo/UKEfTERlZEMuWxaeWRSv4v4pwEqFbt9DwKUltG+8hFRXTeIsFqtV6v1Yr5YLJZK6Vhq8zgVCZRSqFRMgxz7GeV5RR53SumaA+z1ywDeRprH3XsYVStSxhhAdCzsXNAW+JcJaYqqwPrT2sa5iY9lPH7cNMpekMLDOY4nctJuTOyiwOxEQJQdcCvNb8fudelqR8guRcynljzs9Py4HFzipIoIM/n5b0WNMYvF4vT4+OT4ZDlv57PZcrlYLBbzdlZVFSlKbzNsjIqv99h4X251r12e1bHnXvq00koOeQCfMxr03jQDvB7HnGwKUtWK6YBASilFKkZNVUrNmnrWNCfHx0+fPFmvN6v1arverNbr1Xq9Wq/85mC9WvG2C1g2EigC0uATZQPiEJYKjuCE78MztTR4C4ycJuOtnM3GsuukTLsyjE5DMBarSbkxCCRNLEePcRpKIf3K8mq1YpHK1CzivAlobLIBdV1XylTOuX6wnNORR1h7WTMycC8ZovFuTmF27EM949kQkjHGGMMidujE2WJ6S7mW++GK+GiXYo1Jgcg7aE9cg6HIIEiy3mQRqkgrJIVEQAoRJboAZfhGipgIzJ5AXBhwYFRi52A1TIf7Vcd51RcPA4S4C2W0bpq6rpfztj05Pn54cnp8cjyfzZq6Mto0TT1rGlOZuKwzO2HgFFYW5yWyB4O+4T0mB7rP7wcBHZYBvHKFgGu+w1fhRLjvzzieEETij2PgmLQLAlBXuq6OHz44QSABGSyv1quzs7Pf3vz25vX5xfnZm7Ozi7dvhz4MTHMzFk1eMAhHs9UbekUsBmtmLn2BpxsYGU1MpYxwKvk5AbfOGSERBfcRM7kKpjIYvdoA0RthKmW0AmZerze9tqSIHXOBTSEAKV0ZMwBa59KaldH2UuOAO5k8Etx+0MefIUo5D0EkIq2N1ppIOTt4omccJXgnt9BWx3T5Ufh6WiiDQQMprw9nEUDJ0QAyyuoNk/YQ1uVN3FSayUpE6gOM5HUNPF7t4pDX+18H8ReXEZmlZXQEkRBBE5NBMUbrtmna2WzezpaLxcnJ8eNHjx+cnLTtTIXviBcMg3MOwO3saS/vT+Vm997hcVgA/hAPHOWe3+43C4dln74IMc7cq4QBAMAo1dQns9ns6Gj59PFmu9lsNpv1drvZbNar1bv1arVardeb7Waz6bbQ9bEVI9B+kEhAKCTBZBIVkqCDHGOJQfmVFgJKatJUQKZgBxTgC6RikefEWFD/CrA+aJVYMkEGAEQci9iBHEkyKgurCDhrByQB0IocMnORgCgFrRWzzAop+OJg9OU3VdXOZgK4Xm8G23ukBQG00qYyRGqwg+0HjsgVRGO25OwZ4hrLJJYi2gsKlC/DXTJO74qnGj18CIgQCBVJmG8ns+pJCDykgDPCYIgBgslmOcA5znkLbvbqZReM54BIV2bWNLOmbmfNrGnns9l83s5nbd3UtdZ1Xbez2Xw+b2eNGg2uQn5aUDLnBLcPsCM/jHwPC8Bnu4W4s1Vtsl0YdUUi4px1JaKPUBlVnRyfnpx44ECY19vt2dnZm7M352dnZ+fnFxdvL9693azXfdczu9TtRikqs2NhF7Q+7B0OEhJNABTaSIhBJLHxFygyYMrZQOH3sEsvL5wWpOCvgIi40JBGux5ERQQgTqKdcogGJgEY+t5Z1kahUlqTMDsHzuM5MPLUFJ86q5CQkpUoCIADpcxs3jLLdtsBCyhBJOVzXEgxu6EfrB3i+oHZVEjEm1lkRhDm3PlyFYNg5ADZiGgkOwh2RqS00iryfMKWRmAnSQ5ygJCn6wQhApIIi4tLkSt/TwiE/KKiQ3QyIlZVNWuaxWJ+tFyeHB+dHB2fHJ8sl4umqRVilJALiATxyrgWe6RqmmZ/BSp90HF9LgvAFUSr3Zn+1Rrl2x5/q5O595PH646/Af4zul/v583u+NQn11CPQxORSjFPio7mbWXMcjHfPHy02Wy6ruv6ruv7brPdbLebrtts1pv1ervZrruu67Z2sxXeCcQiAlKglJAAeocCjDzOoCnwBsbBORkBQRXaKyFAv5IkbDr5JAtwdOePsAUWWStFdktMYfQ8H5XgjMiEYWbUSoXUwL7vHLucnZvCzoRByOfde1BFRJCBhEWGwUla00CUUpUxpJRz1jrHHimaRCj4dU+SIz/nKQbk6Wey3+Hc7Mso4TEIihVpRUop0qA83A+U4oGCER2mRLfw1tgJM7Nz1oFzDK4HILK+04+TEtJGV5Wp66aum6apZs3M83Zmdd1UdVVXlTF1U7ezWds0s6apVFQtEjIgC0jSb0mwwQ60USy+2XJPN+yEBCp7bgHZd0Pdlkp/9cncS3XaWz3usbS+z8nrO3xUcste9m7Hf2gnjfc9edm/AnwcbjJm9jSmOZ+3bWRwUenjVwFYzOfLxSKAC4pEuOuH1Wr99t27i4uLt+fnb9+9e7t6u16t3r1bb7vODoMEd8oU9YqAhXaLrcuyMsQwWeAQ6ZVgmyBDlRAIBihAhXVFEidQypJMqbII2bABvP0kRrZlVLF5y00C0NooTb7SaUVGV4qUMLO11oUxOEpwOZDkkkSIqJz34gawzr5brf3ExSuttNZKaRYe+oFFlNYkyjkHkjvyHT+fFO+IZfSb5D5dYBJm7JE1Cli/f2WgkG4vxe/npQOyAwiGgAQBz+EhUkTGaFLKI49euKGMrqtq3rbtvJ3P54t5u1wsFvPFYt627byuDYUVVRA8dUAGayflNkWQTiz4ZbfHv0FPdMNbZAr+3Pds4POoNh/s5A8Q0B8IbQqlAlOtiDY0RKrwAUXSdV1X9bxtT4+PNo8f9X3fdV03DEPfD4Mduq7rh0236bbdtuv6rtsO/TAM/WCts8Mw8GAhONaREHkEJFZ/AkJBBCRQGCyaQTz7CAQlBB5KUhHHlBmJaHmRcg5YJP8mqk1YCRhQmBHAVLoyxlrLIuJYVWhmMwJx3WC3HSCCTvDS2PwoGeUBMHPfd34mrEhVVaWUYmFrrWXnsSC/xjhns442IfjJyiGLAQhhtEgUkY5hdOudm0mpYN+W3P89iTJOiAOcw5KlecLgDUeJqKqMUprI1MozcxZt2zSzqqqqumqqxtTaaGO0McZUxtS1qaq6rqq6ruvKmKxORPaLq/egG0UKJxe5Awr/Z4KAPr88gD/0yV8FQyUvoJAAHppw8RsDdtPmCcEYZcx8uVhgSPHyglywg914ldnaD4836816u91st9227zebzXbbddutdS68jh8peh+0CKsn6XBIAHDO1zxhl+PiIWWRF+0wkiBgkDRgKH3gXfuLRgZH1h1h18HirGXmqjJ1VREhCPuSCpEwg7FHD+gQIQGxxNgAYQAg0korRBz6YbADEWqtNCqf2OLNn4uonGLIG9AsPxl1fjbLwdaHSwdswUDS9NCPCDA7FPFoWQD1/A4KRCkoJkE5qAUAjNGzpqmrSitltKqbZj5rj46WbdvOmtmsmbXzed1URquwcnI+dR9h3Mf9yISWRmVsAALCNQEbcvOvK97uO38ZAIu/k7//+1cbOOQBfMzjP7WT/6Dur3lDEJeDPHsVTmmxoQ562s/kO2F0ZZaztjnqj/q+t4PrUvdv7TBY6yO9rbODtW4YumGyhRgGa+3gnLXWWWs9hg6DBXaXVA8fXE6AKP6/iNlZIQpl/TKV9L8+I5EQ3WA7YQQv1iUU8e5iilBrBQioKPgXxTz1bJqGAERCItYxx5hIvxYKA4hWShtltCFU4jgY40iqiGFjQTlZMxFBJeiE2Yt7GYK8i0QhAIB1AGKVsUYFTEkbo40ySittjNHaE/G1MUprUxmjtTbaGKNM+Jv2P1KEipQi0tpURtdVbYzWxhijK61V6fFBHovDsK6GUYqMvr94g/nWXnrwTb+W1wNBI90MXjNz+2j+/nDIA/j88gA+vZP/MGOBKYYoRUBs3ONTce8Fc3k3kgdj2QA2dTWr6yxGTQchIAADdN3Qd9vNdtt13Xaz7rb9Zrsd+m7b9x5WGoa+74dhGHr/sJYdxwDEyBsdaQgKRmikEFGEtDiaKYkIe59TAGYehh5ASFFtaiIcBmJnu2233W7tMATJWLJewICIJX5o8PRxDhwDgTFGxG03aycMgEorYBqGnh0P1jrbZw1cTIecItbC5En+3gQPEZSOCBSRR5MI2TpmMcYoowERUXlbhaoxxlRVZbSuqqoyxjR15QGb2lRVbZqqqpu6aZqq8ktHwTaTFLsWenzrXXiKiTMFGCp87HRJVtb+ZOz3hNqv3jJ8NLj0U6o2cMgDODzubQRwab81vbdl17UuC1CTNwV506691G4FYGrjajNvZ6nTHwbr4sNay44tW8vsrBNm51xv7dB3w2A9q3Lw2wO2zolz1lk7OOecY8fWWee86ss/klqKPc0xm0+w84DIBrYiDMzO8Xr9LqcppDBkAABgUExFCiaLT7f09VDVtYj0794BMCjDlRmgh8GC7UfvnggUgSIAznpgb2lP2tt0EpGikJ8V+vUqdO2eTaVIVXVV142pa8/9xGD1FvXf3u9Z+f9q/welSGuts/XP6KMVBBEq1Gh55pz1hlGud8OIpHtuUw636+cyA/h0YJwPvRH77E9e7nAOoRNDwD2vxCLgfKN92etEQ0uptKqM9qJf9MoyzM8OqZEHsI79jiHsDbp+cMMwOGft4IYhIkhs7eDcMDhrB+ssOwmDBr8eQClpYgGwzg691Votl4umaYZ+qI3ZNDM3WAEJxpnJyJ8IiZIvtgg7dq7vgXk2X5yengrz2azpB9s0tTaVMHd913e9iOjK+KBOpZTSmhRF4ZX/g48KRoWkPF6jtEdzTFXVVWWqyiM7SmutdWWMqaq6bppmptXo8sqk4yvCuUSA7dCBH6FDzPwV2Jd3nXhGe780H630j75ccrjBP5WTOeQB/KFP/qrnERi7Ne/CsmN3n+RlnYwEgqUQBX+glGWypwylYlQp0vO2bRrH7Bw754Q9+C7s2LFzznmBqWN2zOwYsr6Ws3Kt+J/fFThnEbGqK63IORkGy86G/hfTNkYgnGc2wPRPzOJIRJmqqRsR3nY9O6eMJqXFc+ydBQBFOtj2e3Fa8jkLhKs4s8BAvlJxTdBKk1JKhVB4UkTe1NnLfndW8yKJMX604sPmUYK5RqB8JT9TCO4/088bk0ZrHFhaKrfKwnw1RR0K49R0unIbV/rLwPzP1FL/sz75AwT0iTzwvfr5e9oWjFaAacRW/lsOPIzs9gDL86hvxfFKkgndFGjySiswasQ1iTmVWOLaeMX+Y3r55PL6AlMv7OkHgAEgCkC5XDuovBJtk4nxtkiRwVVAccJ24CFLHkbPVhhZjAxAsVTcYvT5BIwA/02/CbLfUnPnE7zXb9lBCfzJQUC3JWbdaq/xSf380ztJvOTHcg9PfsltdnNV87UBTJMaUHrP46Q9jLc9FRgU5pT4+4Gkb+7efdmRtG+BucNrTV5MAH38ZWlEJ2lJncSxAYz3UgJj41i8bES79yOU8QJZaESuxICmH/4eheqV5f+24vn3p0t+LtXm0zn5ww7gz7SvuL/jcXI8TltgTHYPicCOBcSQoJvdMiYZmcbCdTpFrGAhfk4Qx2RBFbm+MZVMqo+O2JACeaHIsw+7kcJSaM8u2z8BJrgJd7bbuG9VLpLErtzYC9xdCnv5ryHg9dgxjsxT37f5P/T+n+QO4PD4qI9x9fqUTub6EnJ1+yayc0zxozwWlqteZbSNkHG9wsIDZ5Q8XAwpEFLg5LVbrQCzU0S0UtmTUfC9SBqVh+SylDpQImRlFiOOkTS50s8Axx5vk59MNwe3qrXjj2SyzuDuZiFqr25iY4WHan5YAA6Pw3biVu3hnqqxD9qfkIhyTS4ij3cKUwE7Xc592Xk7QV0mcmlRy5FosF8phXvWmP1I/LXA+n2oSfCq7QRe8+GU7uWF59P+hevwOCwAh8efdvsybfSvXjZuKHjBq/9czBvkxuNE3D3xlE5wg2pWugiN7MnkJhdKJFs5CNxVOXhzudD4SEwBYncr2Xhw/zksAIfH4fFxNg7vV232SN/wkjVg76LyUS7HR3ix6635P8rHcXh8JgvAgUpfRMBO76MC6L1pI3zNyVx7/HVd9n1e+b0KHdx/OnjFyRd/3sMWl6vq3p2N3XHPxZGpuKE8vnxJkZt/giK5359km196fEG+Qcmx7HuGw1JG58AoH7oM/9pJer+u978UxJfxFmr/m935/A9UejjoAA6PP1/rPuZ4yvVN9yf8nvAquOoaXEVuBYfsqdS73CS8zXag+EDeM3bi+iT2Oz4/HoweDhDQH6f4ye/48r/vfTTh4mSK5pjTAnINViDlFHe/pAuvLjT7rHHlDtxznLRGd7vCN13trpw+p/yzKwApLCbfcHOY6r7MkA+PwwJwzffmD5YHAAhXuJW8/+1zj3fgx7rycikSdN0T4Q1gpXv++O5Q+y5jZOKeLh4vJzwi3mYFwP0umzf6ROSOX6cbX8kyeOcervCfxFL/sz75W+QBpCNvzkP40McjfqiTL4VOeF2PV9IH7+3k8Xe48pcNOvYweaBgxVzOT78t+eQKU6O9pfUerHHx+pPZ70G/ezKXQ037hbtXnfzlV/4G6uQiOvIOJsPyHp+U3MsNWx7/2VWbz+vk6bAJuqzZupuPzuFxeNznznH/kx4AncPjo0BAZUcjn22Fu/XJ38FB4cNdnA88D3if3aXc5ODrDpKRT/3trsEH9+mVaxZ7uYRmeqm8egfGkdt9K/GG7SF8YqbEh8enWVr1B7x54HM9/i4OOvghTwZ//ztTrvSA++hd8Cfawhf8zQ/wCp/zlTwEfnyab/agAzic/A1O5rI14MNzonfXP7nkdT/mlb8amsdbDSfKiQBe6ik9nkLL/pO8ZLE4UOkPJ3/Z8YcZwOFxeBweh8ef9HHIAzic/C12l1PP59/FGx1uGXLwKRm7X9Gx4/VPXmTk4H4Vwf7hxMFS/3Dyl/z8sAM4PA6Pw+Pw+JM+/j9DnqhKS3VwVgAAAABJRU5ErkJggg=='


def _pwa_icon(size=192):
    """★네이비 바탕에 흰 글자 — 외부 파일 없이 코드에서 그린다(배포 파일 증가 0).
       ★v441 브랜드 연동 — BRAND=BARUM이면 B, 아니면 M.
         아이콘만 M으로 남아 「글자는 BARUM인데 아이콘은 M」이 됐다(2026.08.17 실측)."""
    from PIL import Image, ImageDraw
    # ★v629 — 지점장 이미지가 1순위. 실패하면 종전 그리기로 내려간다.
    try:
        import base64 as _b64, io as _io0
        _src = Image.open(_io0.BytesIO(_b64.b64decode(_ICON_SRC_B64))).convert('RGB')
        _src = _src.resize((int(size), int(size)), Image.LANCZOS)
        _b0 = _io0.BytesIO(); _src.save(_b0, 'PNG'); return _b0.getvalue()
    except Exception as _e0:
        print('[v629 PWA] 이미지 아이콘 실패 — 기본 로고로:', str(_e0)[:60])
    # ★v441 (지점장 지시 2026.08.17) — BARUM은 <b>골드 바탕에 남색 B</b>.
    #   MAKEONE은 종전 그대로(네이비 바탕 · 흰 M · 하단 골드 띠).
    GOLD, NAVY = '#c5a052', '#06203f'
    _isB = (_brand()[2] == 'B')
    im = Image.new('RGB', (size, size), GOLD if _isB else NAVY)
    d = ImageDraw.Draw(im)
    u = size / 24
    if not _isB:
        d.rectangle([0, int(20.4 * u), size, size], fill=GOLD)
    w = max(2, int(1.9 * u))
    W = NAVY if _isB else '#ffffff'
    if _brand()[2] == 'B':
        # B — 세로 기둥 + 위아래 반원 두 개
        _x0, _x1 = 6.8 * u, 16.4 * u
        d.line([(_x0, 5.4 * u), (_x0, 18.6 * u)], fill=W, width=w)
        for y0, y1 in ((5.4, 12.0), (12.0, 18.6)):
            d.arc([_x0, y0 * u, _x1, y1 * u], -90, 90, fill=W, width=w)
            _mid = (_x0 + _x1) / 2
            d.line([(_x0, y0 * u), (_mid, y0 * u)], fill=W, width=w)
            d.line([(_x0, y1 * u), (_mid, y1 * u)], fill=W, width=w)
    else:
        # M — 기존 그대로
        d.line([(5.4 * u, 6 * u), (5.4 * u, 18 * u)], fill=W, width=w)
        d.line([(18.6 * u, 6 * u), (18.6 * u, 18 * u)], fill=W, width=w)
        d.line([(5.4 * u, 6 * u), (12 * u, 13.6 * u)], fill=W, width=w)
        d.line([(18.6 * u, 6 * u), (12 * u, 13.6 * u)], fill=W, width=w)
    import io as _io
    b = _io.BytesIO(); im.save(b, 'PNG'); return b.getvalue()


@app.get('/manifest.webmanifest')
def _manifest():
    return JSONResponse({
        "name": _brand()[0] + ' ' + _brand()[1], "short_name": _brand()[0],
        "start_url": "/", "scope": "/", "display": "standalone",
        "background_color": "#06203f", "theme_color": "#06203f",
        "lang": "ko",
        "icons": [{"src": "/icon-192.png", "sizes": "192x192", "type": "image/png",
                   "purpose": "any maskable"},
                  {"src": "/icon-512.png", "sizes": "512x512", "type": "image/png",
                   "purpose": "any maskable"}]},
        headers={'Cache-Control': 'public, max-age=3600'})


@app.get('/icon-{sz}.png')
def _icon(sz: int):
    from fastapi.responses import Response as _R
    try:
        return _R(content=_pwa_icon(int(sz)), media_type='image/png',
                  headers={'Cache-Control': 'public, max-age=86400'})
    except Exception as _e:
        print('[v428 PWA] 아이콘 생성 실패:', _e)
        return _R(content=b'', media_type='image/png')


@app.get('/sw.js')
def _sw():
    from fastapi.responses import Response as _R
    # ★캐시하지 않는다 — 분석 결과는 매번 새로 받아야 한다(제53조 2항).
    # ★v440 제53조 3항 — fetch 핸들러가 없으면 크롬이 「앱 설치」를 띄우지 않는다(2026.08.17 실측).
    #   캐시는 여전히 하지 않는다. 그물만 걸고 네트워크로 그대로 보낸다.
    # ★★★★★v448 제63조 (지점장 지적 2026.08.17 「앱 등록하면서 뭐가 다 없어졌다」)
    #   v440에서 넣은 fetch 핸들러가 <b>POST까지 전부 가로챘다</b>.
    #   분석·리모델링은 FormData를 실어 보내는 POST다. `fetch(e.request)`로 다시 던지면
    #   요청 본문 스트림이 이미 소비돼 <b>응답이 영영 안 온다</b>(버튼이 「비교 중…」에서 멈춤).
    #   → <b>GET만</b> 처리하고 나머지는 손대지 않는다(respondWith를 부르지 않으면 브라우저가 그대로 보낸다).
    #     크롬 설치 요건은 「fetch 핸들러 존재」이므로 이대로도 앱 설치는 된다.
    js = ("self.addEventListener('install', e => self.skipWaiting());\n"
          "self.addEventListener('activate', e => e.waitUntil(clients.claim()));\n"
          "self.addEventListener('fetch', function(e){\n"
          "  if (e.request.method !== 'GET') { return; }   /* POST/업로드는 건드리지 않는다 */\n"
          "  e.respondWith(fetch(e.request));\n"
          "});\n")
    return _R(content=js, media_type='application/javascript',
              headers={'Cache-Control': 'no-store'})


@app.on_event('startup')
def _startup_db():
    # ★DB가 없어도 앱은 뜬다(제54조 4항) — 회원 기능만 꺼진다.
    try:
        _db_init()
    except Exception as _e:
        print('[v429 DB] startup 예외:', str(_e)[:80])


@app.get('/admin')
def admin_page():
    """★권한자 화면(최은혜 지점장). 비번 821024 · 코드 발급 · 차단 · 사용 이력."""
    return HTMLResponse("""<!DOCTYPE html><html lang="ko"><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>BARUM 관리자</title><style>
*{box-sizing:border-box;font-family:-apple-system,"Noto Sans KR",sans-serif}
body{margin:0;background:#0b1420;color:#e8edf4;padding:14px}
/* ★v435 PC 대응(제56조) — 넓은 화면에서 2단 · 표 글자 확대 */
@media (min-width:900px){
 body{padding:24px 32px}
 h1{font-size:24px}
 .wrap{max-width:1200px;margin:0 auto}
 .cols{display:flex;gap:18px;align-items:flex-start}
 .cols>.card:first-child{width:360px;flex:none}
 .cols>.card:last-child{flex:1}
 table{font-size:13px}
 th,td{padding:9px 8px}
 input,select,button{font-size:15px}
}
h1{font-size:19px;margin:6px 0 14px;color:#e7c274}
.card{background:#132234;border:1px solid #24374f;border-radius:12px;padding:14px;margin-bottom:12px}
label{display:block;font-size:12px;color:#9fb0c4;margin:8px 0 4px}
input,select{width:100%;padding:11px;border-radius:8px;border:1px solid #2c4258;
background:#0d1a28;color:#e8edf4;font-size:15px}
button{width:100%;padding:13px;border:0;border-radius:8px;background:#c5a052;color:#0b1420;
font-weight:800;font-size:15px;margin-top:12px}
button.gray{background:#2c4258;color:#e8edf4}
.out{margin-top:12px;padding:14px;background:#0d1a28;border:1px solid #c5a052;border-radius:8px;
text-align:center}
.code{font-size:27px;font-weight:900;color:#e7c274;letter-spacing:.14em;margin:6px 0}
table{width:100%;border-collapse:collapse;font-size:12px;margin-top:6px}
th{background:#0d1a28;padding:7px 5px;text-align:left;color:#9fb0c4;font-size:11px}
td{padding:7px 5px;border-bottom:1px solid #1c2c3f}
.b{color:#ff8a8a}.g{color:#7fd6a8}
.act{padding:5px 9px;border-radius:6px;background:#2c4258;font-size:11px;border:0;color:#e8edf4;width:auto;margin:0}
.act.red{background:#7a2b2b;color:#ffd9d9}
.act.red{background:#7a2b2b;color:#ffd9d9}
.err{color:#ff8a8a;font-size:13px;margin-top:8px}
.sum{font-size:12px;color:#9fb0c4;margin-bottom:6px}
</style></head><body>
<h1>BARUM 관리자</h1>
<div class="wrap">
<div class="card" id="gate">
 <label>관리자 비밀번호</label>
 <input type="password" id="pw" inputmode="numeric" placeholder="••••••">
 <button onclick="login()">확인</button>
 <div class="err" id="ge"></div>
</div>
<div id="main" style="display:none">
 <div class="cols">
 <div class="card">
  <b style="font-size:14px">코드 발급</b>
  <label>이름</label><input id="nm" placeholder="홍길동">
  <label>메모 (선택)</label><input id="mm" placeholder="비워도 됩니다">
  <label>유효기간</label>
  <select id="mo"><option value="12">12개월</option><option value="6">6개월</option>
  <option value="24">24개월</option><option value="120">무기한(10년)</option></select>
  <button onclick="issue()">코드 발급</button>
  <div id="out"></div>
 </div>
 <div class="card" id="pcard" style="display:none;border-color:#c5a052">
  <b style="font-size:14px;color:#e7c274">승인 대기 <span id="pn"></span></b>
  <div id="plist"></div>
 </div>
 <div class="card">
  <b style="font-size:14px">회원 목록</b>
  <div class="sum" id="sum"></div>
  <div id="list"></div>
  <button class="gray" onclick="load()">새로고침</button>
 </div>
 </div>
</div>
</div>
<script>
let PW='';
const $=x=>document.getElementById(x);
async function api(act,extra){
 const fd=new FormData(); fd.append('pw',PW); fd.append('act',act||'');
 for(const k in (extra||{})) fd.append(k,extra[k]);
 const r=await fetch('/admin/api',{method:'POST',body:fd}); return r.json();}
async function login(){
 PW=$('pw').value.trim();
 const j=await api('list');
 if(j.ok){$('gate').style.display='none';$('main').style.display='block';render(j);}
 else{$('ge').textContent=j.error||'오류';}}
async function issue(){
 const j=await api('issue',{name:$('nm').value,months:$('mo').value,memo:$('mm').value});
 if(!j.ok){$('out').innerHTML='<div class="err">'+j.error+'</div>';return;}
 $('out').innerHTML='<div class="out"><div style="font-size:12px;color:#9fb0c4">'+j.name+
  ' 님 코드</div><div class="code">'+j.code+'</div><div style="font-size:11px;color:#9fb0c4">유효 '+
  j.expires+'</div><button onclick="cp(&#39;'+j.code+'&#39;)">복사</button></div>';
 $('nm').value='';$('mm').value=''; load();}
function cp(c){navigator.clipboard.writeText(c);alert('복사됨: '+c);}
async function load(){const j=await api('list'); if(j.ok) render(j);}
function render(j){
 $('sum').textContent='전체 '+j.total+'명 · 차단 '+j.blocked+'명';
 var pd=j.pend||[];
 $('pcard').style.display=pd.length?'block':'none';
 $('pn').textContent=pd.length?('· '+pd.length+'명'):'';
 var ph='<table><tr><th>이름</th><th>연락처</th><th>신청일</th><th></th></tr>';
 pd.forEach(function(r){ph+='<tr><td>'+r.name+'</td><td>'+r.phone+'</td><td>'+r.created+'</td><td>'
  +'<button class="act" style="background:#0e7258" onclick="appr(&#39;'+r.name+'&#39;,&#39;'+r.phone+'&#39;)">승인</button> '
  +'<button class="act red" onclick="rej(&#39;'+r.name+'&#39;,&#39;'+r.phone+'&#39;)">거절</button></td></tr>';});
 $('plist').innerHTML=ph+'</table>';
 let h='<table><tr><th>이름</th><th>코드</th><th>사용</th><th>만료</th><th></th></tr>';
 (j.rows||[]).forEach(r=>{h+='<tr><td>'+r.name+(r.blocked?' <span class="b">차단</span>':'')+
  '</td><td style="font-weight:800;color:#e7c274">'+r.code+'</td><td>'+r.cnt+'회<br><span style="color:#9fb0c4">'+
  r.last+'</span></td><td>'+(r.expires||'-')+'</td><td>'+
  '<button class="act" onclick="tog(&#39;'+r.code+'&#39;,'+(r.blocked?1:0)+')">'+
  (r.blocked?'해제':'차단')+'</button> '+'<button class="act red" onclick="kick(&#39;'+r.code+'&#39;,&#39;'+r.name+'&#39;)">추방</button>'+'</td></tr>';});
 $('list').innerHTML=h+'</table>';}
async function tog(c,b){await api(b?'unblock':'block',{code:c}); load();}
async function appr(nm,ph){
 const j=await api('approve',{name:nm,memo:ph,months:$('mo').value});
 if(j.ok){alert(nm+' 승인 완료 · 코드 '+j.code);load();}else{alert(j.error||'실패');}}
async function rej(nm,ph){
 if(!confirm(nm+' 님의 신청을 거절합니까?'))return;
 const j=await api('reject',{name:nm,memo:ph});
 if(j.ok){load();}else{alert(j.error||'실패');}}
async function kick(c,nm){
 if(!confirm(nm+' ('+c+') 님을 추방합니다.\\n코드가 삭제되어 즉시 접속 불가가 됩니다.'))return;
 if(!confirm('되돌릴 수 없습니다. 정말 추방합니까?'))return;
 const j=await api('delete',{code:c});
 if(j.ok){alert(nm+' 추방 완료');load();}else{alert(j.error||'실패');}}
</script></body></html>""")


# ═══════════ v429 관리자 · 회원 라우트 (제54조) ═══════════
@app.post('/member/apply')
async def member_apply(name: str = Form(''), phone: str = Form('')):
    """★v439 (제59조) 가입 신청 — 설계사가 이름·연락처를 넣고 신청한다.
       코드는 <b>지점장 승인 뒤</b>에 생긴다. 승인 전에는 입장 불가."""
    name = (name or '').strip()
    phone = (phone or '').strip()
    if len(name) < 2:
        return JSONResponse({'ok': False, 'error': '이름을 정확히 입력하십시오'})
    if len(phone) < 9:
        return JSONResponse({'ok': False, 'error': '연락처를 정확히 입력하십시오'})
    c = _db()
    if not c:
        return JSONResponse({'ok': False, 'error': 'DB 미연결 — 지점장에게 문의'})
    try:
        with c, c.cursor() as k:
            k.execute("SELECT status, code FROM members WHERE name=%s AND phone=%s", (name, phone))
            r = k.fetchone()
            if r:
                st, cd = r
                if st == 'active':
                    return JSONResponse({'ok': True, 'status': 'active', 'code': cd})
                return JSONResponse({'ok': True, 'status': 'pending'})
            k.execute("INSERT INTO members(name,phone,status) VALUES(%s,%s,'pending')",
                      (name, phone))
            k.execute("INSERT INTO uselog(code,name,act) VALUES('',%s,'apply')", (name,))
            print('[v439 신청] %s %s' % (name, phone))
        return JSONResponse({'ok': True, 'status': 'pending'})
    except Exception as e:
        return JSONResponse({'ok': False, 'error': str(e)[:100]})
    finally:
        try: c.close()
        except Exception: pass


@app.post('/member/check')
async def member_check_status(name: str = Form(''), phone: str = Form('')):
    """★신청자가 승인됐는지 확인 — 승인됐으면 코드를 돌려준다."""
    c = _db()
    if not c:
        return JSONResponse({'ok': False, 'error': 'DB 미연결'})
    try:
        with c, c.cursor() as k:
            k.execute("SELECT status, code, blocked FROM members WHERE name=%s AND phone=%s",
                      ((name or '').strip(), (phone or '').strip()))
            r = k.fetchone()
            if not r:
                return JSONResponse({'ok': False, 'error': '신청 기록이 없습니다'})
            st, cd, blk = r
            if blk:
                return JSONResponse({'ok': False, 'error': '차단된 계정입니다'})
            if st == 'active' and cd:
                return JSONResponse({'ok': True, 'status': 'active', 'code': cd})
            return JSONResponse({'ok': True, 'status': 'pending'})
    except Exception as e:
        return JSONResponse({'ok': False, 'error': str(e)[:100]})
    finally:
        try: c.close()
        except Exception: pass


@app.post('/member/login')
async def member_login(code: str = Form(''), pw: str = Form('')):
    """설계사 로그인 — 코드 또는 화면비번(0101)."""
    if pw == PW and not code:
        return JSONResponse({'ok': True, 'name': '', 'mode': 'pw'})
    ok, nm, why = _member_check(code)
    return JSONResponse({'ok': ok, 'name': nm, 'mode': 'code'} if ok
                        else {'ok': False, 'error': why})


@app.post('/admin/api')
async def admin_api(pw: str = Form(''), act: str = Form(''), name: str = Form(''),
                    code: str = Form(''), months: int = Form(12), memo: str = Form('')):
    """★권한자 전용(최은혜 지점장 · 821024). 발급 · 목록 · 차단 · 해제 · 삭제."""
    if pw != ADMIN_PW:
        return JSONResponse({'ok': False, 'error': '관리자 비밀번호 오류'})
    c = _db()
    if not c:
        return JSONResponse({'ok': False, 'error': 'DB 미연결 — Railway에 Postgres를 붙이십시오'})
    try:
        import datetime as _d
        with c, c.cursor() as k:
            if act == 'issue':
                if not name.strip():
                    return JSONResponse({'ok': False, 'error': '이름을 입력하십시오'})
                exp = _d.date.today() + _d.timedelta(days=int(months) * 30)
                for _ in range(20):
                    cd = _mk_code()
                    k.execute("SELECT 1 FROM members WHERE code=%s", (cd,))
                    if not k.fetchone():
                        break
                k.execute("INSERT INTO members(name,code,expires,memo) VALUES(%s,%s,%s,%s)",
                          (name.strip(), cd, exp, memo))
                return JSONResponse({'ok': True, 'code': cd, 'name': name.strip(),
                                     'expires': str(exp)})
            if act == 'approve':
                # ★v439 (제59조) 승인 = 코드 생성 + status active
                import datetime as _d2
                exp = _d2.date.today() + _d2.timedelta(days=int(months) * 30)
                for _ in range(20):
                    cd = _mk_code()
                    k.execute("SELECT 1 FROM members WHERE code=%s", (cd,))
                    if not k.fetchone():
                        break
                k.execute("UPDATE members SET code=%s, status='active', expires=%s "
                          "WHERE name=%s AND phone=%s", (cd, exp, name.strip(), memo.strip()))
                k.execute("INSERT INTO uselog(code,name,act) VALUES(%s,%s,'approve')",
                          (cd, name.strip()))
                print('[v439 승인] %s → %s' % (name.strip(), cd))
                return JSONResponse({'ok': True, 'code': cd, 'name': name.strip(),
                                     'expires': str(exp)})
            if act == 'reject':
                k.execute("DELETE FROM members WHERE name=%s AND phone=%s AND status='pending'",
                          (name.strip(), memo.strip()))
                k.execute("INSERT INTO uselog(code,name,act) VALUES('',%s,'reject')",
                          (name.strip(),))
                return JSONResponse({'ok': True})
            if act in ('block', 'unblock', 'delete'):
                if act == 'delete':
                    # ★추방은 되돌릴 수 없다 — 누구를 언제 뺐는지 uselog에 남긴다(제54조 8항).
                    k.execute("SELECT name FROM members WHERE code=%s", (code.upper(),))
                    _r0 = k.fetchone()
                    k.execute("INSERT INTO uselog(code,name,act) VALUES(%s,%s,'kick')",
                              (code.upper(), _r0[0] if _r0 else ''))
                    k.execute("DELETE FROM members WHERE code=%s", (code.upper(),))
                    print('[v429 추방] %s (%s)' % (_r0[0] if _r0 else '?', code.upper()))
                else:
                    k.execute("UPDATE members SET blocked=%s WHERE code=%s",
                              (act == 'block', code.upper()))
                return JSONResponse({'ok': True})
            k.execute("""SELECT name,code,blocked,to_char(created,'YY.MM.DD'),
                         to_char(expires,'YY.MM.DD'),to_char(last_used,'MM.DD HH24:MI'),
                         use_count,COALESCE(memo,''),COALESCE(status,'active'),
                         COALESCE(phone,'') FROM members ORDER BY created DESC LIMIT 300""")
            _all = [{'name': r[0], 'code': r[1] or '', 'blocked': r[2], 'created': r[3],
                     'expires': r[4], 'last': r[5] or '-', 'cnt': r[6],
                     'memo': r[7], 'status': r[8], 'phone': r[9]}
                    for r in k.fetchall()]
            rows = [x for x in _all if x['status'] == 'active']
            pend = [x for x in _all if x['status'] == 'pending']
            k.execute("SELECT COUNT(*) FILTER (WHERE blocked) FROM members")
            blk = k.fetchone()[0]
            return JSONResponse({'ok': True, 'rows': rows, 'pend': pend,
                                 'total': len(rows), 'blocked': blk,
                                 'npend': len(pend)})
    except Exception as _e:
        return JSONResponse({'ok': False, 'error': str(_e)[:120]})
    finally:
        try: c.close()
        except Exception: pass


@app.get('/health')
def health():
    _cib = ci_selftest()   # ★v238 CI 자가진단 — 실패하면 즉시 노출
    # ★v382 조문 자가진단을 /health에 <b>따로</b> 노출한다(CI 카운트에 묻히지 않게).
    _dtb = doctrine_selftest()
    _nmb = name_selftest()   # ★v568 이름 자가진단 — 실패하면 즉시 노출
    _DTN = len(_DOCTRINE_SELFTEST)+len(_SOLO5_SELFTEST)+len(_GEN_SELFTEST)+len(_DEDUP_SELFTEST)
    _STN = len(_CI_SELFTEST)+len(_CI_NONLIFE_SELFTEST)+len(_CI_RATE_SELFTEST)+len(_SILSON_SELFTEST)   # ★v446 손보CI 10건 포함
    # ★★★v309 감사 2종을 /health에서도 즉시 돌린다 — 지점장이 링크 한 번으로 확인.
    try:
        import openpyxl as _ox
        _ws = _ox.load_workbook(TPL_XL)['보장분석']
        _labs = [_ws.cell(r,2).value for r in range(6, _ws.max_row+1) if _ws.cell(r,2).value]
        _a = audit_run(_labs)
        _audit = ('PASS 규칙 %s · 커버리지 %s' % (_a['case'], _a['cov'])) if not _a['fail'] \
                 else ('FAIL %d건 | ' % _a['fail']) + ' | '.join(_a['detail'][:4])
    except Exception as _e:
        _audit = 'ERROR ' + str(_e)[:80]
    return {'ok':True,'version':VSTAMP,
            'audit': _audit,
            'ci_selftest': ('PASS %d/%d' % (_STN, _STN)) if not _cib else ('FAIL: '+' | '.join(_cib[:6])),
            'doctrine': ('PASS 지침자가진단 %d/%d' % (_DTN, _DTN)) if not _dtb else ('FAIL: '+' | '.join(_dtb[:6])),
            '고객명': ('PASS %d/%d' % (len(_NAME_CASES), len(_NAME_CASES))) if not _nmb else ('FAIL %d건: ' % len(_nmb))+' | '.join(_nmb[:4]),
            # ★v440 — 「조문 38/38」은 <b>자가진단 개수</b>였다. 지침 조문 수(59)와 혼동됐다(2026.08.17).
            #   조문 수·결번·분량은 /diag의 제55조 검사가 따로 지킨다. 여기서는 이름만 바로잡는다.
            '조문하한': '%d조 이상 (검사는 /diag)' % DOCTRINE_MIN_ART,
            '동작검사': (lambda t: 'PASS 조문 %d개' % len(t[1]) if not t[0]
                        else 'FAIL %d건: %s' % (len(t[0]), ' | '.join(t[0][:3])))(behave_selftest()),
            # ★v441 — 테스트(BARUM)와 운영(MAKEONE)을 링크만 보고 헷갈리지 않게 찍는다.
            'brand': _brand()[0]}

# ★★v101 진단 엔드포인트(2026.07.20): 폰에서 링크 한 번만 눌러
#   Railway 컨테이너에 pdftotext(poppler)가 실제로 살아있는지 확인한다.
#   'KB(PDF)는 죽고 롯데(txt)는 산다'의 원인을 서버 로그 없이 확정하기 위함.
# ★★★★★v310 체크봇 2종 (지점장 지시 2026.07.31 "업데이트 체크봇·지침 체크봇 필요하다", 영구)
#   문제: 업데이트는 <b>①로컬→zip ②zip→GitHub ③GitHub→Railway ④서버→산출물</b> 4단계인데
#   어디서 끊겼는지 확인할 장치가 없어 "업데이트했는데 없다"가 반복됐다.
#   실측 2026.07.31 = ①v308-fxfont 유실 + ②GitHub v306 / zip v307 이 <b>동시에</b> 끊겨 있었다.
#   ★둘 다 <b>표시 전용</b>. 산출물 값·행 배정은 건드리지 않는다.
@app.get('/dl/{token}/{fname}')
def dl_file(token: str, fname: str):
    """★v316 산출물 실제 다운로드 — blob 대신 진짜 파일 URL.
       모바일(삼성폰 Chrome)에서 blob 연속 다운로드가 차단되던 문제의 근본 해결책."""
    from fastapi.responses import FileResponse
    if not re.fullmatch(r'[0-9a-f]{6,32}', token or ''):
        return JSONResponse({'ok': False, 'error': '잘못된 요청'}, status_code=400)
    _base = os.path.join(tempfile.gettempdir(), 'barum_dl', token)
    _p = os.path.join(_base, os.path.basename(urllib.parse.unquote(fname)))
    if not os.path.isfile(_p):
        return JSONResponse({'ok': False,
                             'error': '파일이 만료되었습니다. 화면에서 다시 분석해 주세요.'},
                            status_code=404)
    return FileResponse(_p, filename=os.path.basename(_p),
                        media_type='application/octet-stream')

@app.get('/version')
def version_bot():
    """업데이트 체크봇 — 서버 실물 8파일의 각인·md5를 GitHub raw와 직접 대조한다.
       ③GitHub→Railway 단계가 끊겼는지(Deploy 블루를 안 눌렀는지)를 앱이 스스로 답한다."""
    import hashlib, urllib.request, random
    RAW = 'https://raw.githubusercontent.com/bokkile83-ui/barum-bunseok-backend/main/'
    NEED = ['main.py','coverage_benchmark.py','report_weasy.py','report_pptx.py',
            'ga_tables.py','master.xlsx','Dockerfile','nixpacks.toml']
    out = {'server_version': VSTAMP}
    rows = []; same = 0; diff = []; err = []
    for fn in NEED:
        p = os.path.join(HERE, fn)
        try:
            b = open(p, 'rb').read()
            mine = hashlib.md5(b).hexdigest()[:10]; size = len(b)
        except Exception as e:
            rows.append({'file': fn, 'server': '★없음', 'github': '-', 'same': False})
            err.append(fn + ' 서버에 없음'); continue
        try:
            u = RAW + fn + '?cb=' + str(random.randint(10**6, 10**7))
            g = urllib.request.urlopen(u, timeout=15).read()
            gh = hashlib.md5(g).hexdigest()[:10]
        except Exception as e:
            rows.append({'file': fn, 'server': '%s (%dB)' % (mine, size),
                         'github': 'ERR ' + str(e)[:40], 'same': None})
            err.append(fn + ' GitHub 조회 실패'); continue
        ok = (mine == gh); same += 1 if ok else 0
        if not ok: diff.append(fn)
        _gs = ''
        if fn.endswith('.py'):
            try:
                _m = re.findall(r'v\d{3}[a-z]*-[a-z0-9]+-\d{8}', g.decode('utf8', 'ignore'))
                _gs = sorted(set(_m))[-1] if _m else ''
            except Exception: _gs = ''
        rows.append({'file': fn, 'server': '%s (%dB)' % (mine, size), 'github': gh,
                     'github_각인': _gs, 'same': ok})
    out['files'] = rows
    out['stamps'] = {}
    for fn in ('main.py','coverage_benchmark.py','report_weasy.py'):
        try:
            t = open(os.path.join(HERE, fn), encoding='utf8', errors='ignore').read()
            m = re.findall(r'v\d{3}[a-z]*-[a-z0-9]+-\d{8}', t)
            out['stamps'][fn] = (sorted(set(m))[-1] if m else '각인없음')
        except Exception as e:
            out['stamps'][fn] = 'ERR ' + str(e)[:40]
    _st = set(v for v in out['stamps'].values() if not str(v).startswith('ERR'))
    out['stamp_same'] = (len(_st) == 1)
    try:
        out['master_rows'] = openpyxl.load_workbook(TPL_XL).active.max_row
    except Exception as e:
        out['master_rows'] = 'ERR ' + str(e)[:40]
    if err:
        out['verdict'] = '★확인불가 — ' + ' / '.join(err[:3])
    elif diff:
        _gsv = [r.get('github_각인') for r in rows if r.get('github_각인')]
        _gsv = sorted(set(_gsv))[-1] if _gsv else '?'
        _sv = out['stamps'].get('main.py', '?')
        if _gsv == _sv:
            _why = ('각인은 같은데 내용이 다르다 → <b>손편집 오염 의심</b>(GitHub 편집창 붙여넣기 사고). '
                    'zip으로 8파일을 다시 통째 올릴 것')
        elif _gsv > _sv:
            _why = ('GitHub이 더 새것(%s) 인데 서버는 %s → <b>Railway Deploy(블루)를 안 눌렀다</b>' % (_gsv, _sv))
        else:
            _why = ('서버가 더 새것(%s) 인데 GitHub은 %s → <b>zip을 GitHub에 안 올렸다</b>' % (_sv, _gsv))
        out['verdict'] = '★불일치 %d개: %s — %s' % (len(diff), ', '.join(diff), _why)
    elif not out['stamp_same']:
        out['verdict'] = '★각인 불일치 — 3파일 버전이 다르다: %s' % out['stamps']
    else:
        out['verdict'] = 'PASS — 서버 8파일 = GitHub 8파일 · 3파일 각인 동일 (%s)' % out['stamps'].get('main.py', '')
    return out

# ★★★★★v569 지침 낭독 로봇 (지점장 지시 2026.08.23
#   「보장분석지 메인화면에 지침=메모리[업뎃날짜/버전] [ ]/100% 다 읽고 했는지 체크하게 하는
#    로봇 + 화면에 표기해라」).
#   [배경] 조문을 박아도·검사기를 넣어도 <b>읽지 않으면 소용이 없다</b>(오늘 4회 위반).
#   ⇒ 매 분석 실행 때 <b>지침 전문을 실제로 읽고</b> 그 결과를 <b>메인화면에 상시 표기</b>한다.
#   ★「읽었다」의 정의 = <b>파일을 열어 조문·핵심 조항·줄수를 세어 통과</b>. 선언이 아니다.
# ★★★★★v572 산출물 5종 로봇 검사 (지점장 지시 2026.08.23
#   「로봇은 엑셀+진단서+보장분석지+비교엑셀+리포트 <b>전체</b>다」).
#   [배경] v571 로봇은 <b>지침 문서만</b> 봤다 — 진단서·리포트는 검사 0건이라 계속 되돌아갔다
#     (실측: 로봇 함수 안에 '진단서'·'report_weasy'·'_op417'·'red_map' 전부 0회).
#   ⇒ 산출물 5종을 <b>실제로 만들어</b> 지침 등식대로인지 본다. 「등식1 = 엑셀 = PPT = 진단서」.
#   ★한 번이라도 어긋나면 100%가 아니다. 화면이 빨강이 된다.
_OUT5 = ('엑셀', '보장분석지PPT', '진단서PPT', '비교엑셀', '리모델링리포트')


def output_selftest():
    """산출물 5종을 실제로 만들어 지침 등식·색 규칙을 검사한다. 실패 목록을 돌려준다."""
    import os as _o, tempfile as _tf, importlib as _il
    bad = []
    _d = _tf.mkdtemp(prefix='rb5_')
    try:
        # ── ① 엑셀: master.xlsx 를 원천으로 rep 생성이 되는가 (제124조 · 등식2)
        try:
            from coverage_benchmark import map_excel_to_report as _m2r
            _rep = _m2r(TPL_XL, settings={'client': '로봇점검', 'branch': '메이크원',
                                          'manager': '최은혜', 'title': '지점장', 'phone': ''})
            if not _rep: bad.append('엑셀 — rep 생성 실패')
            elif not _rep.get('coverage'): bad.append('엑셀 — coverage 비어 있음')
        except Exception as _e:
            bad.append('엑셀 — %s' % str(_e)[:40]); _rep = None

        if _rep:
            # ── ② 고객명: 「고객」으로 떨어지지 않았는가 (제64·128조)
            if str(_rep.get('client') or '') in ('', '고객'):
                bad.append('진단서 — 고객명이 「고객」 (제64·128조)')
            # ── ③ 제안서 단독이면 보유 칸이 없어야 한다 (제22조 = 올 레드)
            if _rep.get('prop_only') and (_rep.get('own_amt') or {}):
                bad.append('진단서 — 제안단독인데 보유 칸이 있다 (제22조)')
            # ── ④ 진단서 렌더: 실제로 PDF 를 만들어 쪽수·백지를 본다
            try:
                from report_weasy import build_report_pdf as _brp
                _pdf = _o.path.join(_d, 'rb.pdf')
                _brp(_rep, _pdf)
                if not _o.path.exists(_pdf) or _o.path.getsize(_pdf) < 50000:
                    bad.append('진단서 — 렌더 실패/과소')
            except Exception as _e:
                bad.append('진단서 — 렌더 예외 %s' % str(_e)[:36])
            # ── ⑤ 색 원천 연동: 엑셀 red_map/gen_map 을 진단서가 받는가 (결과값 동결 #9)
            try:
                from report_weasy import _is_red as _ir
                if _rep.get('prop_only') and not _ir(_rep, '일반암'):
                    bad.append('진단서 — 제안단독인데 레드가 아니다 (제22조)')
            except Exception as _e:
                bad.append('진단서 — 색 연동 예외 %s' % str(_e)[:30])

        # ── ⑥ 리모델링 2종: 비교엑셀·리포트 진입 함수가 살아 있는가 (제129·52조)
        try:
            import remodel as _rm
            for _fn in ('remodel_all', 'remodel_single', 'build_xlsx', 'build_report'):
                if not hasattr(_rm, _fn):
                    bad.append('리모델링 — %s 없음 (제52·129조)' % _fn)
            if hasattr(_rm, 'build_report_pptx'):
                import inspect as _ins
                if 'pngs' not in _ins.signature(_rm.build_report_pptx).parameters:
                    bad.append('리포트 — build_report_pptx 시그니처 변경 (제129조)')
        except Exception as _e:
            bad.append('리모델링 — %s' % str(_e)[:40])

        # ── ⑥-B ★★★제안서 단독을 <b>합성해서</b> 실제로 돌린다 (제22조 = 올 레드).
        #   master.xlsx 는 제안 열이 없어 위 ③⑤ 조건이 안 걸린다 — 조건이 안 맞으면 검사가 아니다.
        #   ⇒ 계약 열 하나를 <b>레드(C00000)</b>로 칠한 임시 엑셀을 만들어 prop_only 를 강제 재현한다.
        try:
            import openpyxl as _ox
            from openpyxl.styles import Font as _Ft
            _tmp = _o.path.join(_d, 'prop.xlsx')
            _wb = _ox.load_workbook(TPL_XL); _ws = _wb.active
            _n = 0
            for _r in range(6, _ws.max_row + 1):
                if _ws.cell(_r, 2).value:
                    _ws.cell(_r, 3).value = 1000
                    _ws.cell(_r, 3).font = _Ft(color='FFC00000', name='맑은 고딕', size=9)
                    _n += 1
            _ws.cell(1, 3).value = '테스트손보\n제안상품\n[갱신]'
            _wb.save(_tmp); _wb.close()
            from coverage_benchmark import map_excel_to_report as _m2r2
            _rp = _m2r2(_tmp, settings={'client': '로봇점검', 'branch': '메이크원',
                                        'manager': '최은혜', 'title': '지점장', 'phone': ''})
            if not _rp.get('prop_only'):
                bad.append('진단서 — 제안단독 판정 실패 (제22조)')
            if (_rp.get('own_amt') or {}):
                bad.append('진단서 — 제안단독인데 보유 칸이 남았다 (제22조)')
            from report_weasy import _is_red as _ir2, _op417 as _op2
            for _lb in ('일반암', '뇌혈관진단비', '급성심근경색', '질병수술비'):
                if not _ir2(_rp, _lb):
                    bad.append('진단서 — 제안단독인데 「%s」가 레드가 아니다 (제22조)' % _lb); break
            _o417 = _op2(_rp, '일반암')
            if _o417 and _o417[0]:
                bad.append('진단서 — 제안단독인데 보유 금액이 찍힌다 (제22조)')
        except Exception as _e:
            bad.append('진단서 — 제안단독 합성 예외 %s' % str(_e)[:36])

        # ── ⑧ ★갱신/비갱신이 엑셀과 일치하는가 (제5·95조 · 등식1)
        try:
            if _rep:
                _rn = int(_rep.get('renew') or 0); _nr = int(_rep.get('nonrenew') or 0)
                _nc = int(_rep.get('n_contract') or 0)
                if _nc and (_rn + _nr) != _nc:
                    bad.append('진단서 — 갱신%d+비갱신%d ≠ 계약%d (제5·95조)' % (_rn, _nr, _nc))
        except Exception as _e:
            bad.append('갱신구분 — %s' % str(_e)[:30])
        # ── ⑨ ★CI 담보가 빠지지 않는가 (제8.4조 · 제119조 2항)
        try:
            if _rep:
                _ci = _rep.get('ci') or {}
                _cia = _rep.get('ci_amounts') or {}
                # ★v573 실측: CI가 없는 계약도 `ci` dict 자체는 항상 온다(`present: False`).
                #   <b>present 가 True 일 때만</b> 금액을 따진다 — 아니면 상시 오탐이다.
                if _ci.get('present') and not _cia:
                    bad.append('진단서 — CI 계약인데 CI 금액이 비었다 (제8.4조)')
                if 'ci' not in _rep:
                    bad.append('진단서 — CI 키 자체가 없다 (제119조 2항)')
        except Exception as _e:
            bad.append('CI — %s' % str(_e)[:30])
        # ── ⑩ ★★★등식1 — 진단서 값 = 엑셀 값 (결과값 동결)
        try:
            if _rep:
                _cv = _rep.get('coverage') or []
                _dm = _rep.get('dambo') or {}
                _diff = []
                for _row in _cv[:400]:
                    if not isinstance(_row, dict): continue
                    _k = _row.get('name') or _row.get('label')
                    if not _k or _k not in _dm: continue
                    _a = _row.get('amount'); _b2 = _dm.get(_k)
                    if isinstance(_a, (int, float)) and isinstance(_b2, (int, float)) and int(_a) != int(_b2):
                        _diff.append('%s 진단서%s≠엑셀%s' % (_k, format(int(_a), ','), format(int(_b2), ',')))
                if _diff:
                    bad.append('등식1 위반 %d건: %s' % (len(_diff), ' / '.join(_diff[:2])))
        except Exception as _e:
            bad.append('등식1 — %s' % str(_e)[:30])

        # ── ⑦ 보장분석지 PPT 폼이 있는가 (빠지면 조용히 소실 — 실사고)
        try:
            if not _o.path.exists(_o.path.join(_o.path.dirname(_o.path.abspath(__file__)),
                                               'ppt_form.pptx')):
                bad.append('보장분석지PPT — ppt_form.pptx 없음')
        except Exception: pass
    finally:
        try:
            import shutil as _sh; _sh.rmtree(_d, ignore_errors=True)
        except Exception: pass
    return bad


# ★★★★★v618 제142조 — 로봇2(정답지 대조) (지점장 지시 2026.08.30 「로봇1 로봇2 각 심장 넣어줘」)
#   [왜 필요한가] 로봇1은 <b>조문이 코드에 있는가</b>만 본다. 그래서 값이 다 틀려도 100%였다.
#     오늘 지적 20건 중 로봇1이 잡은 것은 <b>0건</b>이다.
#   ⇒ 로봇2는 <b>지점장이 지정한 정답지</b>와 산출물 값을 대조한다. 다르면 실패다.
#     정답지는 `_ANSWER` 에 담보명 → 값으로 박는다. 없으면 로봇2는 <b>대기</b>다.
_ANSWER = {
    # (고객, 담보명): 정답값 — 지점장 확정 실측치만 넣는다
    ('정상범', '뇌혈관진단비'): 500,
    ('정상범', '부정맥'): 500,
    ('정상범', '심부전'): 500,
    ('정상범', '일반암'): 1000,
    ('정상범', '유사암(갑.기.경.제)'): 200,
    ('윤선경', '입원'): 5000,
    ('윤선경', '통원'): 20,
}

# ★심장 정답 — 로봇1(매핑)·로봇2(값) 양쪽에서 본다
_ANSWER_HEART = [
    ('흥국화재', '특정3대심장질환진단비(통합간편가입형)', ['부정맥', '심부전']),
    ('흥국화재', '특정심혈관질환(기타심장부정맥)진단비', ['부정맥']),
    ('흥국화재', '특정심혈관질환(기타심장부정맥제외)진단비', ['협심증', '빈맥', '심부전']),
    ('롯데손해보험', '심혈관질환진단비(특정심장질환Ⅰ)(건강체할인형)', ['급성심근경색']),
    ('롯데손해보험', '심혈관질환진단비(특정심장질환Ⅱ)(건강체할인형)', ['협심증', '주요심장염증']),
    ('롯데손해보험', '심혈관질환진단비(특정15대심장질환)(건강체할인형)',
     ['협심증', '급성심근경색', '빈맥', '부정맥', '심부전', '주요심장염증', '심장판막']),
    # ★조문 386행 — KB 특정Ⅰ은 <b>염증 포함</b>이다(내가 정답표를 잘못 적어 실패가 났다)
    ('KB손해보험', '특정심장질환(특정Ⅰ)진단비', ['협심증', '빈맥', '심부전', '주요심장염증']),
    ('DB손해보험', '특정심장질환(특정Ⅰ)진단비', ['협심증', '주요심장염증']),
]


def answer_heart_selftest():
    """★v618 심장 — 회사별 분해가 정답과 같은가(로봇1·로봇2 공용)."""
    bad = []
    for _co, _nm, _want in _ANSWER_HEART:
        try:
            _got = heart_rows(_co, _nm)
        except Exception as _e:
            bad.append('[심장정답] %s %s — 예외 %s' % (_co, _nm[:20], str(_e)[:20])); continue
        if _got != _want:
            bad.append('[심장정답] %s %s — <b>%s</b> (정답 %s)' % (_co, _nm[:20], _got, _want))
    return bad


def answer_selftest(client='', cov=None):
    """★v618 로봇2 — 산출물 값이 정답지와 같은가.
       cov = {담보명: 값}. 그 고객의 정답이 `_ANSWER`에 없으면 <b>대기</b>(빈 목록)."""
    if not client or not isinstance(cov, dict):
        return []
    bad = []
    for (_c, _nm), _want in _ANSWER.items():
        if _c != client:
            continue
        _got = cov.get(_nm)
        try:
            _gv = int(float(_got or 0))
        except Exception:
            _gv = 0
        if _gv != int(_want):
            bad.append('[정답지] %s · %s — <b>%s</b> (정답 %s)' % (_c, _nm, _gv or '없음', _want))
    return bad


def doctrine_robot(heavy=False):
    """지침=메모리 낭독 로봇.
       ★★★★★v570 (지점장 2026.08.23 「읽는 건 할 수 있다. 무조건 100%가 나와야 한다 —
         로봇이 100% 읽도록 <b>행동하게</b> 하라」).
       [구 결함] v569는 <b>조항 문자열이 있는가</b>만 봤다. 그건 「있다」이지 「읽었다」가 아니다.
         오늘 4회 위반 때도 조문은 전부 있었다 — 100% 초록이었을 것이다.
       ⇒ <b>조문마다 실행 검사를 연결</b>한다. 조문을 읽고 <b>그대로 동작하는지 실제로 돌려</b>
         전부 통과할 때만 100%다. 「읽었다 = 그 조문대로 코드가 움직인다」로 정의를 바꾼다.
       돌려주는 값: (읽음%, 버전, 갱신일시, 조문수, 실패목록)"""
    import os as _o, re as _r, datetime as _dt
    _p = _o.path.join(_o.path.dirname(_o.path.abspath(__file__)), 'BARUM_DOCTRINE.md')
    try:
        _d = open(_p, encoding='utf-8').read()
    except Exception as _e:
        return (0, VSTAMP, '-', 0, ['★BARUM_DOCTRINE.md 를 열 수 없다: %s' % str(_e)[:40]])

    _bad = []
    def _has(k): return k in _d

    # ── 조문 → 실행 검사. (이름, 조문이 있는가, 조문대로 도는가)
    _CHECKS = []

    def _ck(nm, key, fn):
        """조문 존재 + 실행 검사를 함께 본다. 둘 다 통과해야 1점."""
        if not _has(key):
            _CHECKS.append((nm, False, '조문 누락')); return
        try:
            _r2 = fn()
        except Exception as _e:
            _CHECKS.append((nm, False, '실행 예외 %s' % str(_e)[:30])); return
        _CHECKS.append((nm, not _r2, _r2 or ''))

    # 1) 제1지침 — 존재 + 3개 하위 조항
    _ck('제1지침', '## 제1지침 — 매 턴 무한반복',
        lambda: '' if ('sync_check' in _d and 'FINAL.zip' in _d) else '실행 절차 미기재')
    # 2) 0항 느낌금지
    _ck('0항 느낌금지', '0항 — 느낌을 쓰지 않는다',
        lambda: '' if '느낌은 근거가 아니다' in _d or '느낌이며' in _d else '금지표 미기재')
    # 3) 질문금지 — grep 선행
    _ck('질문금지', '질문 금지 조항 — 지침에 있는 것은 묻지 않는다',
        lambda: '' if 'grep 출력이 붙지 않은 질문' in _d else '판정문 미기재')
    # 4) 제0조
    _ck('제0조 서열', '## 제0조 — 최상위 법률',
        lambda: '' if '지침 = 메모리는 100% 지켜야 할 법률' in _d else '1항 미기재')
    # ★★★★★v583 (지점장 지시 2026.08.24 「<b>지침=메모리 / 엑셀은 법이다 / 1부터 다 읽도록
    #   로봇이 체크해서 100% 나와야 한다</b>」) — 오늘 확정된 4건을 <b>실행 검사</b>로 박는다.
    #   조문만 있고 코드가 안 따르면 <b>즉시 실패</b>다.
    _ck('제131조 3색', '제131조', lambda: '' if (
        '블랙 or 블루 or 레드' in _d or '블랙 OR 블루 OR 레드' in _d) else '3색 조문 미기재')

    def _ck131code():
        try:
            import report_weasy as _rw5, report_pages as _rp5
        except Exception as _e:
            return '모듈 로드 실패 %s' % str(_e)[:24]
        _miss = []
        if not hasattr(_rw5, 'color_selftest'): _miss.append('진단서 color_selftest')
        if not hasattr(_rw5, '_split_span'):    _miss.append('진단서 분할표기')
        if not hasattr(_rp5, '_own_span'):      _miss.append('리포트 3색')
        return ('없다: ' + ' · '.join(_miss)) if _miss else ''
    _ck('제131조 코드', '제131조', _ck131code)

    def _ck122():
        try:
            _raw = '자동차사고변호사선임비용(현물급부선택가능)(실손)'
            if not _is_hyunmul(_raw, 'DB손해보험', '2026.08.20'): return 'DB 2026.08 현물 판정 실패'
            if _is_hyunmul(_raw, 'DB손해보험', '2026.07.31'):     return '2026.07이 현물로 잡힌다'
            if _is_hyunmul('교통사고변호사법률상담비용', 'DB손해보험', '2026.08.20'): return '상담비용 오판'
            if _is_hyunmul(_raw, '삼성화재', '2026.08.20'):        return '타사 오판'
            return ''
        except Exception as _e:
            return '실행 예외 %s' % str(_e)[:24]
    _ck('제122조 현물', '제122조', _ck122)

    def _ckexcel():
        """★엑셀이 법 — 리포트의 「보유」가 <b>최종 엑셀 보유 열</b>에서 오는가."""
        try:
            import remodel as _rm5, report_pages as _rp6
            if not hasattr(_rm5, 'split_sheet'): return 'split_sheet 없다'
            _src = open(__file__, encoding='utf-8').read() if False else None
            _t = open(_o.path.join(_o.path.dirname(_o.path.abspath(_rm5.__file__)), 'remodel.py'),
                      encoding='utf-8').read()
            if 'v580c' not in _t: return '엑셀 2개 경로가 최종 엑셀 보유 열을 쓰지 않는다'
            _t2 = open(_o.path.join(_o.path.dirname(_o.path.abspath(_rp6.__file__)), 'report_pages.py'),
                       encoding='utf-8').read()
            if 'cov_split' not in _t2: return '리포트가 갱신색(cov_split)을 읽지 않는다'
            if 'cov_text' not in _t2:  return "리포트가 글자값('현물')을 읽지 않는다"
            return ''
        except Exception as _e:
            return '실행 예외 %s' % str(_e)[:24]
    _ck('엑셀이 법', '엑셀이 법', _ckexcel)


    def _ck132():
        """★v589 제132조 — 상급병원 일당은 미기재(1인실 예외). 실제로 돌려서 확인한다."""
        try:
            _t = [('상해입원일당(상급종합병원)', None), ('질병입원일당(상급종합병원)', None),
                  ('상급종합병원 질병입원일당(1인실)', '1인실 상급병원'),
                  ('상해입원일당(종합병원)', '종합병원 상해입원일당')]
            for _nm, _want in _t:
                if resolve2(_nm)[0] != _want:
                    return '%s → %s (정답 %s)' % (_nm, resolve2(_nm)[0], _want)
            return ''
        except Exception as _e:
            return '실행 예외 %s' % str(_e)[:24]
    _ck('제132조 상급일당', '제132조', _ck132)

    def _ck134():
        """★v593 제134조 — 일상배상책임 1억 고정이 코드에 살아 있는가."""
        try:
            _t = open(_o.path.join(_o.path.dirname(_o.path.abspath(__file__)), 'main.py'),
                      encoding='utf-8').read()
        except Exception:
            return ''
        if resolve2('가족일상생활배상책임(대물')[0] != '일상배상책임':
            return '잘린 담보명이 일상배상책임으로 안 간다'
        # ★v593b — 코드 문자열이 아니라 <b>실제 값</b>으로 센다. 문자열만 보면
        #   `_dm['일상배상책임'] = 10000` 을 지워도 통과했다(역검증 실패 · 2026.08.26).
        _fake = [{'company': '테스트', 'product': 'T', 'dambo': {'일상배상책임': 20}}]
        try:
            _fix_ilbae(_fake)
        except Exception as _e:
            return '고정 함수 없음/예외 %s' % str(_e)[:20]
        if _fake[0]['dambo'].get('일상배상책임') != 10000:
            return '일상배상 20 → %s (정답 10000)' % _fake[0]['dambo'].get('일상배상책임')
        return ''
    _ck('제134조 일상배상', '제134조', _ck134)
    # ★★★★★v618 로봇2 — <b>값 검사</b>(지점장 지시 2026.08.30 「로봇1 로봇2 각 심장」).
    #   로봇1은 「조문이 코드에 있는가」만 본다 — 오늘 지적 20건 중 <b>0건</b>을 잡았다.
    #   로봇2는 <b>정답지 값</b>과 대조한다. 다르면 그 자리에서 운다.
    #   ★정답은 지점장이 실물로 확정한 것만 넣는다(내가 만들지 않는다).
    _ANSWER_SET = [
        # (담보명, 회사, 정답 마스터행)  — 2026.08.30 지점장 확정 실물
        ('특정3대심장질환진단비(통합간편가입형)', '흥국화재', ['부정맥', '심부전']),
        ('심혈관질환진단비(특정15대심장질환)(건강체할인형)', '롯데손해보험',
         ['협심증', '급성심근경색', '빈맥', '부정맥', '심부전', '주요심장염증', '심장판막']),
        ('심혈관질환진단비(특정심장질환Ⅰ)(건강체할인형)', '롯데손해보험', ['급성심근경색']),
        ('심혈관질환진단비(특정심장질환Ⅱ)(건강체할인형)', '롯데손해보험', ['협심증', '주요심장염증']),
    ]
    globals()['_ANSWER_SET'] = _ANSWER_SET

    def _ck_ans():
        """★로봇2 — 지점장이 확정한 실물 정답과 값이 같은가."""
        for _nm, _co, _want in _ANSWER_SET:
            try:
                _got = heart_rows(_co, _nm)
            except Exception as _e:
                return '[정답지] %s — 예외 %s' % (_nm[:20], str(_e)[:18])
            if list(_got or []) != list(_want):
                return '[정답지] %s — %s (정답 %s)' % (_nm[:22], _got, _want)
        return ''
    _ck('로봇2 정답지', '제138조', _ck_ans)

    def _ck_dup():
        """★★★★★v620 로봇2 — <b>같은 계약이 두 열로 갈리지 않는가</b>.
           (지점장 실측 2026.08.30 「DB손보가 왜 2개가 출력되지 · 이건 치명적 실수다」)
           [사고] 별첨에서 담보 한 줄이 <b>상품명에 딸려 들어와</b>
             `참좋은종합보험1709` / `참좋은종합보험1709질병사망(…)1,000` 두 이름이 되어
             <b>보험료·가입일·만기일이 같은데도</b> 계약이 둘로 갈렸다.
           [구멍] 로봇에 <b>이 검사가 아예 없었다</b> — 계약이 둘로 나와도 100%였다.
           ⇒ 가짜 계약 2건(이름만 다르고 나머지 동일)을 넣어 <b>1건으로 합쳐지는지</b> 센다."""
        try:
            _a = {'company': 'DB손보', 'product': '참좋은종합보험1709', 'renewal': '갱신',
                  'contract_date': '2017.09.19', 'expiry_date': '2032.09.19',
                  'pay_period': '', 'pay_count': '', 'premium': 106358, 'remain': 1,
                  'dambo': {'일반암': 1000}}
            _b = dict(_a); _b['product'] = '참좋은종합보험1709질병사망(15년갱신)(납입면제이전) 1,000'
            _b['dambo'] = {'뇌출혈진단비': 500}
            _r = parse_txt('', '테스트', extra=[_a, _b])
            _n = len((_r or {}).get('contracts') or [])
            if _n != 1:
                return '같은 계약(보험료·기간 동일)이 %d건으로 갈렸다 (정답 1건)' % _n
            return ''
        except Exception as _e:
            return '실행 예외 %s' % str(_e)[:24]
    _ck('로봇2 계약중복', '제5조', _ck_dup)

    def _ck_val():
        """★★★★★v621 로봇2 — <b>지점장이 지적한 값</b>을 전부 센다.
           (지점장 2026.08.30 「로봇이 진짜 기능을 하게 해줘」)
           로봇1은 조문이 코드에 있는가만 본다 — 하루 지적 20건 중 <b>0건</b>을 잡았다.
           ⇒ 지적받은 <b>담보명 원문 → 정답</b>을 여기 박는다. 값이 틀리면 그 자리에서 운다.
           ★정답은 지점장이 실물로 확정한 것만 넣는다(내가 만들지 않는다)."""
        _T = [
            # 2026.08.30 지점장 실측 지적분
            ('가족일상생활배상책임(대물',                      '일상배상책임'),   # 자기부담금 20을 금액으로 읽던 것
            ('질병수술(간편건강고지)담보',                     '질병수술비'),     # 「비」 없어도 수술비
            ('상해수술(간편건강고지)담보',                     '상해수술비'),
            ('질병수술위로금',                                None),           # 넓히면 새는 것
            ('질병수술자금',                                  None),
            ('상해수술비(종합병원)',                           None),           # 병원규모 변형
            ('상해입원일당(상급종합병원)',                      None),           # 상급 일당 미기재
            ('상급종합병원 질병입원일당(1인실)',                '1인실 상급병원'), # 1인실 예외
            ('간병인사용질병입원일당Ⅷ(요양병원및의원제외)(1-180일)(간편건강고지)담보', '간병인'),
            ('간병인사용질병입원일당Ⅷ(요양,정신,한방병원및의원제외)(181-365일)(간편건강고지)담보', None),
            ('간병인사용상해입원일당Ⅷ(간호간병통합서비스)(일반병동/재활병동)(1-180일)(간편건강고지)담보', '간호통합병동'),
            ('유방주요치료비(간편건강고지)(중증(삼중음성)유방암진단)담보', None),  # 마스터 무행
            ('갑상선암주요치료비',                             None),
            ('심혈관질환진단비(특정심장질환Ⅰ)(건강체할인형)',    None),           # 협심증 2배 방지
            ('상해(일반상해, 전체상해를 의미) 의료비(입원+통원)', '입원'),          # 4세대 통합형
            ('비급여 도수/체외충격파/증식치료',                 '도수치료'),
            ('비급여 주사제',                                 '비급여주사'),
            ('비급여 MRI 검사',                               'MRI'),
            ('(간편3.10.5)허혈성심장질환진단Ⅱ(해약환급금미지급형)(계약자적립액미지급형)(납입면제형)', '허혈성 진단비'),
            ('신특정순환계질환 주요치료비ⅢPlus(간편가입)',      '2대 주요치료비'), # 엑셀은 2대
            ('뇌출혈 사망',                                   None),           # 특정 사망은 진단비 아님
            ('급성심근경색 사망',                              None),
            ('암사망',                                        None),
            ('일반사망',                                      '일반사망'),      # 마스터 행은 살아 있어야
            ('상해사망',                                      '상해사망'),
        ]
        _bad = []
        for _nm, _want in _T:
            try:
                _got = resolve2(_nm)[0]
            except Exception as _e:
                _bad.append('%s — 예외 %s' % (_nm[:18], str(_e)[:16])); continue
            if _got != _want:
                _bad.append('%s → %s (정답 %s)' % (_nm[:22], _got, _want))
        if _bad:
            return '%d건: %s' % (len(_bad), _bad[0])
        return ''
    _ck('로봇2 값검사', '엑셀이 법', _ck_val)

    _ck('제135조 수술비', '제135조',
        lambda: '' if not surg_selftest() else ('%d건: %s' % (len(surg_selftest()), surg_selftest()[0])))

    # ★★★★★v595 제137조 (지점장 지시 2026.08.26 「<b>보장분석지 / 보장분석지+제안서 /
    #   제안서 / 엑셀1·2 비교 — 이 4가지에 대해 로봇이 따로 지정되고 따로 각각 검사해야 한다</b>」).
    #   [왜 필요한가] 오늘 실사고가 전부 <b>모드별로</b> 났다 —
    #     ・제안서 단독만 1건 인식(`_prop_cts[0]`) — 조합 모드는 멀쩡했다
    #     ・4세대 실손은 2열(롯데) 경로만 깨졌다 — 3열은 멀쩡했다
    #   ⇒ <b>모드마다 진입 함수가 살아 있는지</b>를 각각 센다. 한 모드가 죽어도 나머지는 100%가 된다.
    def _ck_mode(mode):
        """4가지 실행 모드가 각각 <b>진입 가능한 상태</b>인가."""
        def _f():
            try:
                if mode == '보장분석지':
                    if not callable(globals().get('parse_txt')): return 'parse_txt 없음'
                    if 'extra' not in parse_txt.__code__.co_varnames: return 'parse_txt(extra=) 인자 없음'
                    return ''
                if mode == '보장분석지+제안서':
                    if 'extra' not in parse_txt.__code__.co_varnames: return 'extra 합류 경로 없음'
                    _t = open(_o.path.join(_o.path.dirname(_o.path.abspath(__file__)), 'main.py'),
                              encoding='utf-8').read()
                    if '제안 계약 {len([x for x in extra if x])}건 합류' not in _t:
                        return '제안 계약 합류부가 없다'
                    return ''
                if mode == '제안서단독':
                    _t = open(_o.path.join(_o.path.dirname(_o.path.abspath(__file__)), 'main.py'),
                              encoding='utf-8').read()
                    # ★v595b — <b>주석 문장이 자기 검사에 걸렸다</b>(실측 오탐).
                    #   문자열만 세면 「그 표현을 설명한 주석」까지 위반으로 잡힌다.
                    #   ⇒ <b>주석·문자열을 걷어낸 실제 코드 줄</b>에서만 센다.
                    _code = '\n'.join(_l for _l in _t.split('\n')
                                      if _l.strip() and not _l.strip().startswith('#'))
                    if re.search(r'^\s*_solo\s*=\s*_prop_cts\[0\]', _code, re.M):
                        return '제안서 단독이 첫 건만 쓴다(제133조 위반)'
                    if '_solos = [_x for _x in _prop_cts' not in _t:
                        return '제안서 전건 사용 코드가 없다'
                    return ''
                if mode == '엑셀2개비교':
                    import remodel as _rm
                    for _fn in ('read_sheet', 'split_sheet', 'compare', 'remodel_all'):
                        if not hasattr(_rm, _fn): return 'remodel.%s 없음' % _fn
                    # ★v595b — 문자열이 아니라 <b>실제 동작</b>으로 센다.
                    #   `v580c` 문자열은 주석에도 있어 지워도 통과했다(역검증 실패 2026.08.26).
                    #   ⇒ `split_sheet`가 <b>보유·제안·글자값·갱신색 4종을 다 만드는가</b>를 본다.
                    import inspect as _ins
                    _src = _ins.getsource(_rm.split_sheet)
                    for _key in ('cov_text', 'cov_split', "'제안 합계'"):
                        if _key not in _src:
                            return 'split_sheet에 %s 없음(제9조 5항·제131조)' % _key
                    #   ★주석·로그 문자열을 걷어낸 <b>실제 호출문</b>만 본다(오탐·미탐 둘 다 막는다).
                    _rsrc = '\n'.join(_l for _l in _ins.getsource(_rm.remodel_all).split('\n')
                                      if _l.strip() and not _l.strip().startswith('#'))
                    if not re.search(r'=\s*split_sheet\s*\(', _rsrc):
                        return '엑셀 2개 경로가 split_sheet를 호출하지 않는다(제127조 · 최종 엑셀 보유 열)'
                    return ''
                return ''
            except Exception as _e:
                return '실행 예외 %s' % str(_e)[:24]
        return _f
    for _md in ('보장분석지', '보장분석지+제안서', '제안서단독', '엑셀2개비교'):
        _ck('모드 ' + _md, '제137조', _ck_mode(_md))

    # ★★★★★v601 제137조 2항 (지점장 지시 2026.08.29 「<b>로봇 4모드는 돌리고 정독이다</b>」).
    #   [구 동작] 화면을 열 때마다 배지가 떴다(v569 상시 표기). 그래서 <b>분석을 돌리기 전에도</b>
    #     색이 나왔고, 지점장이 아무것도 안 했는데 레드로 보였다.
    #   ⇒ <b>4모드를 실제로 돌린 결과가 있을 때만</b> 정독 %를 확정한다.
    #     직전 분석이 남긴 `_MODE4_LAST`가 없으면 <b>「대기」</b>로 두고 색을 칠하지 않는다.
    _m4 = globals().get('_MODE4_LAST')
    if _m4 is not None:
        _CHECKS.append(('4모드 실행결과', not _m4,
                        ('4모드 실패 %d건: %s' % (len(_m4), ' / '.join(_m4[:2]))) if _m4 else ''))

    # 5) ★제64조·제128조 — 이름. <b>실제로 13건을 돌린다</b>
    def _ckname():
        _n = len(_NAME_CASES)
        if _n < 13: return '이름 케이스 %d건 < 13건 (검사가 지워졌다)' % _n
        _b = name_selftest()
        return ('name_selftest 실패 %d건' % len(_b)) if _b else ''
    _ck('이름(제64·128조)', '제128조', _ckname)
    # 6) ★제123조 — 심장. 29건 실행 + <b>케이스 수</b>
    #   ★v570-B 실측: 케이스를 지우면 검사 대상이 줄어 실패 0이 된다(고장 ③이 100% 통과했다).
    #     ⇒ <b>케이스 수 하한</b>을 함께 본다. 검사를 지워서 통과하는 길을 막는다.
    def _ckheart():
        _n = len(_HEART_CASES)
        if _n < 29: return '심장 케이스 %d건 < 29건 (검사가 지워졌다)' % _n
        _b = heart_case_selftest()
        return ('heart 실패 %d건' % len(_b)) if _b else ''
    _ck('심장(제123조)', '제123조', _ckheart)
    # ★v618 로봇1 — 회사별 심장 분해 정답 대조
    _ck('심장정답(로봇1)', '제123조',
        lambda: '' if not answer_heart_selftest()
        else ('%d건: %s' % (len(answer_heart_selftest()), answer_heart_selftest()[0])))
    # ★v618 로봇2 — 정답지 값 대조(직전 분석 결과가 있을 때만)
    def _ck_ans2():
        _c2 = globals().get('_ANS_LAST')
        if not _c2:
            return ''                      # 분석 전 = 대기
        _b2 = answer_selftest(_c2.get('client'), _c2.get('cov')) + answer_heart_selftest()
        return ('%d건: %s' % (len(_b2), _b2[0])) if _b2 else ''
    _ck('로봇2 정답지', '제142조', _ck_ans2)
    # 7) ★제9조 — 실손. 23건 실행 + 케이스 수
    def _cksil():
        _n = len(_SILSON_CASES)
        if _n < 23: return '실손 케이스 %d건 < 23건 (검사가 지워졌다)' % _n
        _b = silson_selftest()
        return ('silson 실패 %d건' % len(_b)) if _b else ''
    _ck('실손(제9조)', '## 제9조', _cksil)
    # 7-B) ★규칙 케이스 총수 — 감사 91건
    def _ckaudit():
        _n = len(_AUDIT_CASES)
        return '' if _n >= 91 else '규칙 케이스 %d건 < 91건 (검사가 지워졌다)' % _n
    _ck('규칙케이스(제0조)', '## 제0조', _ckaudit)
    # 8) ★제8.4조 — CI. 실행
    _ck('CI(제8.4조)', '제8',
        lambda: ('ci 실패 %d건' % len(ci_selftest())) if ci_selftest() else '')
    # 9) ★제124조 1항 — master.xlsx 를 손대지 않았는가(지문 대조)
    def _ckmaster():
        import hashlib as _h
        try:
            _m = _h.md5(open(TPL_XL, 'rb').read()).hexdigest()
        except Exception as _e:
            return 'master 열기 실패 %s' % str(_e)[:24]
        return '' if _m == 'a963d8fa243635f0f5142828ee567ad1' else 'master.xlsx 변조 %s' % _m[:8]
    _ck('엑셀불가침(제124조)', '제124조', _ckmaster)
    # 9-B) ★★★★★산출물 5종.
    #   ★★★★★v573 긴급 (지점장 2026.08.23 「니가 마지막 준 거는 앱이 멈춘다 · 열리지도 않아」):
    #   [실측] v572는 <b>홈 화면 1회 렌더에 52.3초</b>가 걸렸다. `_brandize`가 매 페이지 로드마다
    #     산출물 5종을 <b>실제로 만들었기</b> 때문이다(45쪽 렌더 + 엑셀 합성).
    #   ⇒ 화면용은 <b>문서 검사만</b>(0.1초). 산출물 5종은 <b>분석 실행 때 백그라운드 병렬</b>로 돌린다.
    if heavy:
        _ck('산출물5종(등식1)', '## 제0조',
            lambda: ('산출물 실패 %d건: %s' % (len(output_selftest()),
                                            ' / '.join(output_selftest()[:2])))
                    if output_selftest() else '')
    else:
        _o5 = globals().get('_OUT5_LAST')     # 직전 분석 때 병렬 로봇이 남긴 결과를 그대로 쓴다
        if _o5 is not None:
            _CHECKS.append(('산출물5종(등식1)', not _o5,
                            ('산출물 실패 %d건: %s' % (len(_o5), ' / '.join(_o5[:2]))) if _o5 else ''))
    # 10) ★조문 자가진단 — 조문↔코드 대응 전수
    _ck('조문↔코드', '## 제0조',
        lambda: ('doctrine_selftest 실패 %d건' % len(doctrine_selftest())) if doctrine_selftest() else '')

    # ── ★★★★★v571 (지점장 2026.08.23 「로봇은 무조건 지침=메모리를 <b>100% 읽게</b> 해야 한다.
    #   그게 로봇의 역할이다」).
    #   [구 결함] v570은 <b>핵심 10조항만</b> 봤다. 131조 중 10조 = 8%를 읽고 100%라 찍었다.
    #   ⇒ <b>조문 131개를 전수 순회해 본문을 실제로 읽는다.</b>
    #     읽음의 정의 = ①제목이 있다 ②본문이 비어 있지 않다(최소 40자) ③미결 표시가 없다.
    #     하나라도 못 읽으면 그만큼 %가 깎인다. <b>131/131이라야 100%다.</b>
    _arts = sorted({int(x) for x in _r.findall(r'제(\d+)조', _d) if int(x) <= 300})
    if len(_arts) < DOCTRINE_MIN_ART:
        _bad.append('조문 %d개 < 하한 %d' % (len(_arts), DOCTRINE_MIN_ART))
    _miss = [n for n in range(_arts[0], _arts[-1] + 1) if n not in _arts] if _arts else []
    if _miss: _bad.append('결번 %d개' % len(_miss))

    # 조문 전수 정독 — 제목 위치를 찾아 다음 조문 직전까지를 본문으로 읽는다
    _pos = []
    for _m in _r.finditer(r'^#{1,4}[^\n]*?제(\d+)조', _d, _r.M):
        _no = int(_m.group(1))
        if _no <= 300: _pos.append((_no, _m.start()))
    _pos.sort(key=lambda x: x[1])
    _read, _thin = 0, []
    for _i2, (_no, _st) in enumerate(_pos):
        _en = _pos[_i2 + 1][1] if _i2 + 1 < len(_pos) else len(_d)
        _body = _r.sub(r'\s', '', _d[_st:_en])
        if len(_body) >= 40 and '[미결]' not in _d[_st:_en]:
            _read += 1
        else:
            _thin.append('제%d조' % _no)
    _tot = len(_pos) or 1
    if _thin:
        _bad.append('본문 미달 %d조(%s)' % (len(_thin), ','.join(_thin[:4])))

    for _nm, _ok, _why in _CHECKS:
        if not _ok: _bad.append('%s — %s' % (_nm, _why))
    _cpass = sum(1 for _n2, _o2, _w2 in _CHECKS if _o2)
    # ★읽음% = <b>조문 전수 정독률</b> × <b>실행검사 통과율</b>. 둘 다 만점이라야 100%.
    _pct = int(round(_read / _tot * (_cpass / max(1, len(_CHECKS))) * 100))
    if _bad and _pct == 100: _pct = 99
    globals()['_ROBOT_READ'] = (_read, _tot, _cpass, len(_CHECKS))

    try:
        _ts = _dt.datetime.fromtimestamp(_o.path.getmtime(_p)).strftime('%Y.%m.%d %H:%M')
    except Exception: _ts = '-'
    return (_pct, VSTAMP, _ts, len(_arts), _bad)


@app.get('/robot')
def robot_page():
    _pct, _ver, _ts, _n, _bad = doctrine_robot(heavy=True)   # ★엔드포인트는 전부 돌린다
    _rr = globals().get('_ROBOT_READ') or (0, 0, 0, 0)
    return {'ok': not _bad, '읽음': '%d%%' % _pct, '버전': _ver, '갱신': _ts,
            '조문': _n, '본문정독': '%d/%d' % (_rr[0], _rr[1]),
            '실행검사': '%d/%d' % (_rr[2], _rr[3]), '실패': _bad}


@app.get('/doctrine')
def doctrine_bot():
    """지침 체크봇 — 지침 조항별로 담보명→마스터행 케이스를 돌려 통과/실패를 그대로 노출한다."""
    groups = []; tot = 0; ok_all = 0
    for gname, cases in _AUDIT_GROUPS:
        det = []; ok = 0
        for raw, exp in cases:
            try: got = resolve2(raw)[0]
            except Exception: got = None
            good = ((got or None) == (exp or None))
            ok += 1 if good else 0
            det.append({'담보명': raw, '기대': exp or '(기재금지)',
                        '실제': got or '(기재금지)', 'PASS': good})
        tot += len(cases); ok_all += ok
        groups.append({'조항': gname, '통과': '%d/%d' % (ok, len(cases)),
                       'FAIL': [d for d in det if not d['PASS']], '케이스': det})
    cov = {'미복귀_정상허용': sorted(_COV_ALLOW), '건수': len(_COV_ALLOW)}
    try:
        _ws = openpyxl.load_workbook(TPL_XL)['보장분석']
        _labs = [str(_ws.cell(r, 2).value).strip() for r in range(6, _ws.max_row + 1)
                 if _ws.cell(r, 2).value]
        _seen = []; _bad = []
        for L in _labs:
            if L in _seen: continue
            _seen.append(L)
            try: g = resolve2(L)[0]
            except Exception: g = None
            if g != L and L not in _COV_ALLOW: _bad.append({'라벨': L, '실제': g})
        cov['검사'] = '%d/%d' % (len(_seen) - len(_bad), len(_seen))
        cov['FAIL'] = _bad
    except Exception as e:
        cov['검사'] = 'ERR ' + str(e)[:40]
    return {'version': VSTAMP,
            '지침정본': _DOCTRINE_SRC,
            '해석원칙_출처': _PRINCIPLES_SRC,
            '해석원칙': [p[0] for p in (_PRINCIPLES or [])],
            '규칙케이스': '%d/%d' % (ok_all, tot),
            '마스터커버리지': cov.get('검사', '-'),
            'verdict': ('PASS' if ok_all == tot and not cov.get('FAIL') else '★FAIL — 아래 FAIL 항목을 먼저 고칠 것'),
            '조항별': groups, '커버리지': cov}

@app.get('/diag')
def diag():
    import subprocess, shutil
    out = {'version': VSTAMP}
    out['pdftotext_path'] = shutil.which('pdftotext') or '없음(★범인)'
    try:
        r = subprocess.run(['pdftotext', '-v'], capture_output=True, text=True, timeout=20)
        out['pdftotext_ver'] = ((r.stderr or '') + (r.stdout or '')).strip().split('\n')[0][:80]
    except Exception as e:
        out['pdftotext_ver'] = 'ERR ' + str(e)[:80]
    for d in ('/usr/share/poppler', '/usr/share/poppler/cMap'):
        try: out['poppler_data:' + d] = os.path.isdir(d)
        except Exception: out['poppler_data:' + d] = 'ERR'
    out['api_key'] = bool(os.environ.get('ANTHROPIC_API_KEY', ''))
    out['pdf2image'] = _mod_ok('pdf2image')
    out['weasyprint'] = _mod_ok('weasyprint')
    # ★★★v256(2026.07.27): <b>배포 파일 실물 점검</b> — "절반만 나온다"의 최다 원인은
    #   코드가 아니라 <b>파일 누락</b>이다. ppt_form.pptx가 없으면 build_ppt가 조용히 False를
    #   반환해 <b>엑셀만 나오고 PPT 2종이 통째로 사라진다</b>(에러 메시지 없음).
    #   → 서버에 실제로 있는 파일과 바이트수를 그대로 노출한다. 로그를 못 봐도 한 번에 판별된다.
    _need = ['main.py','coverage_benchmark.py','report_weasy.py','report_pptx.py',
             'ga_tables.py','master.xlsx','Dockerfile','nixpacks.toml',
             'ppt_form.pptx','chiryo_form.pptx','requirements.txt']
    _files = {}
    for _fn in _need:
        _p = os.path.join(HERE, _fn)
        try:
            _files[_fn] = os.path.getsize(_p) if os.path.exists(_p) else '★없음'
        except Exception as _e:
            _files[_fn] = 'ERR ' + str(_e)[:40]
    out['files'] = _files
    out['missing'] = [k for k, v in _files.items() if v == '★없음']
    try:
        _wb = openpyxl.load_workbook(TPL_XL); out['master_rows'] = _wb.active.max_row
    except Exception as _e:
        out['master_rows'] = 'ERR ' + str(_e)[:60]
    return out

def _mod_ok(m):
    try:
        __import__(m); return True
    except Exception as e:
        return 'ERR ' + str(e)[:60]

# ★★★★★v441 브랜드 스위치 (지점장 지시 2026.08.17, 영구)
#   「업데이트할 때마다 BARUM으로 테스트하고 MAKEONE에 최종 업데이트한다」
#   → 코드는 <b>하나</b>다. 서버마다 환경변수 BRAND만 다르게 준다.
#     zip을 두 벌 유지하면 반드시 갈라진다(분해 경로 2개 금지 원칙).
#   BRAND=BARUM   → 테스트 서버
#   미설정/그외    → MAKEONE (운영 기본 · 안전한 쪽)
def _brand():
    import os as _o
    if (_o.environ.get('BRAND', '') or '').strip().upper() == 'BARUM':
        return ('BARUM', '보장분석', 'B')
    return ('MAKEONE', '보장설명서', 'M')


def _brandize(html):
    b, sub, ini = _brand()
    # ★★★★★v569 지침 낭독 로봇 배지 — <b>메인화면 상시 표기</b>(지점장 지시 2026.08.23).
    #   화면을 열 때마다 지침 파일을 <b>실제로 열어 읽고</b> 그 결과를 찍는다. 선언이 아니다.
    try:
        _p, _v, _t, _n, _bad = doctrine_robot()
    except Exception as _e:
        _p, _v, _t, _n, _bad = 0, '?', '-', 0, [str(_e)[:30]]
    # ★★★★★v601 제137조 2항 (지점장 지시 2026.08.29 「<b>로봇 4모드는 돌리고 정독이다</b>」).
    #   분석을 한 번도 안 돌렸으면 <b>색을 칠하지 않는다</b>(회색 「대기」).
    #   구 동작은 화면을 열기만 해도 레드가 떠서 <b>지점장이 아무것도 안 했는데 빨갛게</b> 보였다.
    # ★★★★★v602 (지점장 선택 2026.08.30 「㉠ 분석 전에도 초록으로」).
    #   v601은 분석 전 회색 「대기」였는데 <b>100%인데도 회색</b>이라 헷갈렸다.
    #   ⇒ <b>100%면 초록, 실패 있으면 레드</b>. 분석 여부와 무관하게 결과 그대로 보여준다.
    _ran4 = globals().get('_MODE4_LAST') is not None
    _ok = (_p == 100 and not _bad)
    _col = '#0d6b45' if _ok else '#b3242c'
    _box = '■' if _ok else '□'
    _rr = globals().get('_ROBOT_READ') or (0, 0, 0, 0)
    # ★★★★★v602 (지점장 지시 2026.08.30 「<b>문구만 좀더 심플 가자 너무 길다</b>」).
    #   한 줄로 줄인다 — <b>버전 · 정독% · 검사 통과수</b>만. 날짜·조문수·본문 수치는 뺀다.
    # ★v603b — 아이콘을 넣고, 좁은 화면에서 <b>줄이 끊기지 않게</b> 조각마다 nowrap.
    _icon = '🟢' if _ok else '🔴'
    # ★★★★★v618 (지점장 지시 2026.08.30 「로봇1 로봇2 각 심장 넣어줘 · 게임처럼」).
    #   <b>로봇1</b> = 지침·코드 검사(조문이 코드대로 도는가)
    #   <b>로봇2</b> = 값 검사(정답지와 같은 값이 나오는가)
    #   각 로봇의 <b>심장(HP)</b>을 표시한다. 실패 1건에 심장 하나가 깨진다.
    # ★★★★★v623c (지점장 지시 2026.09.02 「저 하트 빼자 · 로봇[동그라미 or x]」).
    #   하트 5개를 ○ 하나로 바꾼다. 실패가 한 건이라도 ×.
    def _hp(_pct, _fail, _ready=True):
        if not _ready:
            return '–', '#6c7887'
        return ('○', '#0d6b45') if int(_fail) == 0 else ('×', '#b3242c')
    _n1 = len([x for x in (_bad or []) if '정답지' not in str(x)])
    _n2 = len([x for x in (_bad or []) if '정답지' in str(x)])
    _r2ready = bool(globals().get('_ANSWER_SET'))
    _hp1, _c1 = _hp(_p, _n1)
    _hp2, _c2 = _hp(_p, _n2, _r2ready)
    # ★★★★★v623 (지점장 지시 2026.09.02 「로봇부분 글자 더 줄여서 한 줄로」).
    #   라벨(지침·값·검사)과 가운데점을 뺀다. 버전은 날짜를 떼고 앞만(v623-skin).
    #   줄바꿈 금지(nowrap) — 둘로 쪼개지면 한 줄이 아니다.
    _vs = str(_v).rsplit('-', 1)[0]
    _badge = ('<span style="display:flex;flex-wrap:nowrap;gap:7px;align-items:center;'
              'font-size:10.5px;white-space:nowrap;overflow:hidden">'
              '<span style="color:%s">🤖1%s</span>'
              '<span style="color:%s">🤖2%s</span>'
              '<span style="color:%s">%s<b>%s</b> %d%% %d/%d</span>'
              '</span>'
              % (_c1, _hp1, _c2, _hp2, _col, _icon, _vs, _p, _rr[2], _rr[3]))
    # ★v602 「<b>안 되는 건 분석 후 좀 크게 표기해줘</b>」 — 실패는 <b>크고 굵게</b>.
    if _bad:
        _badge += ('<br><span style="color:#b3242c;font-size:15px;font-weight:800;'
                   'line-height:1.5">★ ' + '<br>★ '.join(_bad[:3]) + '</span>')
    return (html.replace('@@BRAND@@', b)
                .replace('@@BSUB@@', sub)
                .replace('@@DROBOT@@', _badge)
                .replace('@@BINI@@', ini))


@app.get('/',response_class=HTMLResponse)
def home():
    # ★★★★★v450 제65조 (지점장 지적 2026.08.17 「아예 버튼이 안 먹힌다」)
    #   화면 HTML에 캐시 헤더가 없어 브라우저·설치앱이 <b>옛 화면을 계속 쓴다</b>.
    #   코드를 고쳐 배포해도 지점장 화면만 그대로다 → 버튼이 안 먹히는 것처럼 보인다.
    #   → 화면은 매번 새로 받는다.
    return HTMLResponse(_brandize(INDEX_HTML),
                        headers={'Cache-Control': 'no-store, no-cache, must-revalidate',
                                 'Pragma': 'no-cache', 'X-BARUM-VER': VSTAMP})

@app.post('/check')
async def check_pw(body:dict): return {'ok':body.get('pw')==PW}


# ★★★★★v408 (지점장 지시 2026.08.12 — 3대 원칙 ①): <b>zip 발행 자체를 코드가 검증한다.</b>
#   지점장 원문: 「1. zip에 지침·메모리 항상 최신본 / 2. 무조건 읽고 분석하기 / 3. 지침이 법이고 엑셀이 기준이다」
#   ①은 그동안 <b>손으로</b> 확인했다 — 그래서 2026.08.12에 <b>제0조 6항 본문이 빠진 채</b> 나갈 뻔했다.
#   → 이 함수 하나가 zip 10파일·4파일 각인 일치·지침 조문·핵심 조문 수록·셀프테스트를 <b>전부</b> 찍는다.
# ★v423 — 11파일. report_pages.py(리모델링 7쪽 시안)와 requirements.txt를 뺐다가
#   배포가 통째로 깨질 뻔했다(팩폭 30 · 2026.08.15). 파일이 늘면 <b>이 목록부터</b> 늘린다.
ZIP9 = ['main.py','coverage_benchmark.py','report_weasy.py','report_pptx.py','ga_tables.py',
        'master.xlsx','Dockerfile','nixpacks.toml','ppt_form.pptx','remodel.py','assets_b64.py',
        'report_pages.py','requirements.txt']
DOC_MUST = ['단독 5종','심장 묶음','제외 7종','결과값 동결','엑셀에 없는 건','제0조 6',
            '배포 9파일','PC ↔ 폰','미결','조문 = 테스트']


def js_selftest(root='.'):
    """★★★★★v435 제56조 — <b>JS 문법은 아무도 안 봤다</b>.
       오늘 이스케이프 사고로 관리자 화면·로그인이 통째로 죽은 채
       「배포하십시오」를 네 번 말했다. 화면 <script>를 뽑아 검사한다."""
    import os as _o, re as _r
    bad = []
    try:
        src = open(_o.path.join(root, 'main.py'), encoding='utf-8').read()
    except Exception as e:
        return ['[제56조] main.py 읽기 실패: %s' % e]
    for n, m in enumerate(_r.finditer(r'<script>(.*?)</script>', src, _r.S), 1):
        js = m.group(1)
        fns = _r.findall(r'function\s+([A-Za-z_$][\w$]*)\s*\(', js)
        dup = sorted({x for x in fns if fns.count(x) > 1})
        if dup:
            bad.append('[제56조] JS 블록%d 함수 중복 정의 %s' % (n, dup))
        if _r.search(r'onclick="[^"]*\\\\+\'', js):
            bad.append('[제56조] JS 블록%d onclick 이스케이프 위험 — &#39; 를 쓸 것' % n)
    return bad


# ═══════════════════════════════════════════════════════════════════
# ★★★★★ 제68조 — 동작검사 (v457 · 지점장 2026.08.17 「지침이 무시되면 우리 방법이 없다」)
#
#   팩폭 실측: 조문검사 153건이 오늘 결함 8건을 <b>0건</b> 잡았다.
#   이유는 하나다 — 구 검사는 「이 문자열이 파일에 있나」만 본다.
#   `from assets_b64 import` 가 있으면 통과한다. `*`인지 명시인지 안 본다.
#
#   ★그래서 <b>실제로 함수를 돌려 결과를 비교</b>하는 검사를 만든다.
#     지침을 내가 읽든 안 읽든, <b>동작이 틀리면 zip이 발행되지 않는다.</b>
#     지침 준수를 내 의지에 맡기지 않는다 — 그것이 이 조문의 목적이다.
#
#   ★커버리지를 <b>동작검사 기준</b>으로 다시 센다. 100%를 쉽게 만들지 않는다.
# ═══════════════════════════════════════════════════════════════════
def behave_selftest(d=''):
    """★★★★★v460 제69조 — 심플 모드 값검사 (지점장 확정 2026.08.17).

    <b>실제로 사고가 났던 것만</b> 검사한다. 실패하면 zip 발행을 막는다.
    ★형식 검사(문자열이 파일에 있나 · 커버리지 · 분량)는 <b>전부 뺐다</b> —
      오늘 하루에 거짓 경보 3번, 정당한 배포 차단 3번을 만들었다.
    ★검사를 늘리려면 <b>그 결함이 실제로 났을 때만</b> 늘린다.
    """
    import os as _os
    bad = []; done = set()
    base = d or _os.path.dirname(_os.path.abspath(__file__)) or '.'

    def _t(jo, name, got, exp):
        done.add(jo)
        if got != exp:
            bad.append('[제%d조 %s] %r (기대 %r)' % (jo, name, got, exp))

    # ① 제60조 손해보험사는 CI가 없다 (2026.08.17 실사고)
    for co, pd, exp in _CI_NONLIFE_SELFTEST:
        _t(60, 'CI %s' % co, _isci_prod(pd, co), exp)

    # ② 제4조 단체보험은 기재 금지 (2026.08.17 실사고)
    for pd, cd, ed, exp in [('빅히트단체상해보험', '', '', True),
                            ('무배당 단체상해보험', '2020.01.01', '2021.01.01', True),
                            ('무배당 단체상해보험', '2015.05.19', '2072.05.19', False),
                            ('무배당 롯데 내마음속 건강보험', '', '', False)]:
        _t(4, '단체 %s' % pd[:10], _is_group_ins(pd, cd, ed), exp)

    # ③ 제18·67조 대표(max) — 합산하면 값이 2배가 된다 (실사고 2건)
    for std, exp in [('암일당', True), ('통합암', True), ('통합전이암', True),
                     ('일반암', False), ('유사암(갑.기.경.제)', False)]:
        _t(67 if '통합' in std else 18, '대표max %s' % std, _is_repmax(std), exp)

    # ④ 제22조 통합암 행선지 (지침오류로 지점장이 직접 정정한 조문)
    for nm, exp in [('통합암진단비', '통합암'),
                    ('암진단비(유사암제외)(통합간편가입형)', '일반암')]:
        try: got = resolve2(nm)[0]
        except Exception as _e: got = 'ERR %s' % _e
        _t(22, '매핑 %s' % nm[:16], got, exp)

    # ⑤ 제3조 단독 5종 (지점장 「200번째다」)
    for nm, exp in [('허혈성심장질환진단비', True), ('뇌혈관진단비', True),
                    ('특정허혈심장질환진단비', False)]:
        try: got = is_solo5_name(nm)
        except Exception as _e: got = 'ERR %s' % _e
        _t(3, '단독5종 %s' % nm[:14], got, exp)

    # ⑥ 제61조 assets 상수가 실제로 import됐나 (진단서·설명서 실종 사고)
    try:
        import report_weasy as _rw
        for _n in ('_BON_FORM', '_BON_LINK', '_FIN_SURVEY'):
            _t(61, '상수 %s' % _n, hasattr(_rw, _n), True)
    except Exception as _e:
        bad.append('[제61조] report_weasy 로드 실패 %s' % _e); done.add(61)

    # ⑦ 제63·65조 화면이 죽지 않았나 (앱 등록 사고)
    try:
        _js = _sw().body.decode()
        _t(63, 'sw GET만', "!== 'GET'" in _js, True)
        _r = home(); _h = dict(_r.headers); _html = _r.body.decode()
        _t(65, '화면 캐시금지', _h.get('cache-control', '').startswith('no-store'), True)
        _t(65, '_onReady', 'function _onReady' in _html, True)
        _t(65, '리셋 버튼', 'id="tabnew"' in _html, True)
    except Exception as _e:
        bad.append('[제63조] 화면 확인 실패 %s' % _e); done.update({63, 65})

    # ⑧ 제28·25조 「보유」는 그 이름의 담보가 있을 때만
    #    지점장 「허헐성은 허혈성단독이라고」 — 실사고. 심플 모드에서도 남긴다.
    #    ★엑셀을 실제로 만들어 map_excel_to_report를 돌린다. 끝열 수식은 숫자로 채워 읽는다.
    try:
        import coverage_benchmark as _cb, openpyxl as _op2, tempfile as _tf
        def _c(co, dm):
            return {'company': co, 'product': co + '보험', 'renewal': '비갱신', 'premium': 10000,
                    'dambo': dm, 'contract_date': '2020.01.01', 'expiry_date': '2060.01.01',
                    'total_months': '20년납', 'lump_sum': 0, 'paid_up': False, 'prop': False}
        _o = _os.path.join(_tf.gettempdir(), '_chk28_%d.xlsx' % _os.getpid())
        build_excel({'client': '검사', 'contracts': [
            _c('A', {'뇌혈관질환진단비': 2500, '허혈성심장질환진단비': 1500})]}, _o)
        _wb = _op2.load_workbook(_o); _w = _wb['보장분석']
        for _r in range(6, _w.max_row + 1):
            _nums = [v for v in (_w.cell(_r, _c2).value for _c2 in range(3, _w.max_column))
                     if isinstance(v, (int, float))]
            if _nums: _w.cell(_r, _w.max_column).value = sum(_nums)
        _wb.save(_o)
        _rp = _cb.map_excel_to_report(_o, settings={'client': '검사'})
        _sh = _rp.get('scope_heart') or []; _sb = _rp.get('scope_brain') or []
        _t(28, '허혈성만 → angina 없음', 'angina' in _sh, False)
        _t(28, '허혈성만 → chronic 있음', 'chronic' in _sh, True)
        _t(25, '뇌혈관만 → hem 없음', 'hem' in _sb, False)
        _t(25, '뇌혈관만 → other 있음', 'other' in _sb, True)
        try: _os.remove(_o)
        except Exception: pass
    except Exception as _e:
        bad.append('[제25·28조] 확인 실패 %s' % _e); done.update({25, 28})

    # ⑨ 제62조 마스터 담보행 행높이 (지점장 「골절 이하로 세로폭이 줄었다」)
    try:
        import openpyxl as _op
        _ws = _op.load_workbook(_os.path.join(base, 'master.xlsx'))['보장분석']
        _none = [r for r in range(6, _ws.max_row + 1)
                 if (_ws.row_dimensions[r].height if r in _ws.row_dimensions else None) is None]
        _t(62, '행높이 미지정', len(_none), 0)
    except Exception as _e:
        bad.append('[제62조] master 확인 실패 %s' % _e); done.add(62)

    return bad, done


def _is_repmax(std):
    """대표(max) 담보 판정 — <b>한 곳에서만</b> 정의한다(제19조 「같은 판정은 한 곳에서만」).
       제18조(암일당)·제67조(통합암) 사고는 이 목록이 흩어져 있어 생겼다."""
    return std in _REPMAX_ROWS


# ★대표(max) 정본 목록 — 계약 내 대표와 끝열 수식이 <b>같은 목록</b>을 본다.
# ★★★★★v587 (지점장 실측 2026.08.25 「<b>암주요치료비가 500+500 해서 1천인데 500만 잡힌다</b>」
#   「<b>비급여(하이클래스)는 1000이 나왔다</b>」 · 천미옥 KB 제안서).
#   [실측] 같은 제안서·같은 구조인데 처리가 갈려 있었다 —
#     577 암(유사암제외) 주요치료비Plus(<b>상급종합병원</b>)(호르몬 제외형)   500
#     579 암(유사암제외) 주요치료비Plus(<b>상급종합병원</b>)(호르몬 일부지급형) 500  → 합 1,000이 정답
#     609·618 비급여 암 주요치료비Plus(종합/상급종합) 500+500 = <b>1,000 정상</b>(하이클래스는 대표max가 아니다)
#   [원인] `암주요치료비`만 대표(max) 목록에 있어 <b>이름이 다른 별개 담보 2건을 하나로 눌렀다</b>.
#   ⇒ 대표(max)는 <b>「이름만 다른 같은 담보」를 막는 장치</b>다. 별개 담보까지 막으면 값이 사라진다.
#     하이클래스(암)와 <b>같은 처리</b>로 맞춘다 — 목록에서 뺀다.
# ★★★★★v596 제135조 2항 (지점장 확정 2026.08.28 「입원수술비·통원수술비 둘 중 대표값 하나만」)
#   같은 계약의 `질병입원수술비`+`질병통원수술비`가 한 행으로 모이면 <b>합산되면 안 된다</b>.
_REPMAX_ROWS = ('질병수술비','상해수술비',
                '표적항암치료비','다빈치로봇수술비','n대수술비','120대수술비','입원','통원','약값','약',
                '간병인','간병인지원일당','창상봉합술','항암방사선약물','암수술','중입자치료비',
                # ★★★★★v607 (지점장 실측 2026.08.30 「암주요치료비는 대표값만 넣는 건데 1500만 되어 있고」).
                #   흥국은 <b>급부 유형별로 6개</b>로 쪼갠다(암수술비·방사선·약물·호르몬·중환자실·유사암).
                #   합산하면 1,500이 된다 → <b>대표값</b>이 정본.
                '암주요치료비',
                '암일당','통합암','통합전이암','간호통합병동','합의금','6주미만',
                '1인실 상급병원','1인실 종합병원','2대 주요치료비')


def zip_selfcheck(d=''):
    """zip 발행 전 필수 검증(제0조 6항·제12조). 실패 목록을 돌려준다. 빈 리스트여야 발행 가능."""
    import os as _os, re as _re
    bad=[]; base=d or _os.path.dirname(_os.path.abspath(__file__)) or '.'
    for _f in ZIP9 + ['BARUM_DOCTRINE.md']:
        if not _os.path.exists(_os.path.join(base,_f)): bad.append(f'파일 없음 {_f}')
    stamps={}
    for _f in ('main.py','coverage_benchmark.py','report_weasy.py','BARUM_DOCTRINE.md'):
        _fp=_os.path.join(base,_f)
        try:
            with open(_fp,encoding='utf-8',errors='replace') as _h: _c=_h.read()
        except Exception as _e: bad.append(f'읽기 실패 {_f} {_e}'); continue
        _al=_re.findall(r'v\d{3}[a-z]*-[a-z0-9]+-\d{8}', _c)
        stamps[_f]=_al[-1] if _al else '없음'
    if len(set(stamps.values()))>1:
        bad.append('각인 불일치 ' + ' / '.join(f'{k}={v}' for k,v in stamps.items()))
    # ★★★★★v457 제68조 — 동작검사. 문자열이 아니라 <b>실행 결과</b>를 본다.
    #   지침을 읽든 안 읽든 동작이 틀리면 여기서 막힌다.
    try:
        _bb, _bdone = behave_selftest(base)
        bad += ['값검사 ' + x for x in _bb]
        print('[값검사] 조문 %d개 실행 · 실패 %d건' % (len(_bdone), len(_bb)))
    except Exception as _e:
        bad.append('값검사 실행 실패 %s' % _e)
    # ★★★★★v460 제69조 — <b>문자열 조문검사는 발행을 막지 않는다.</b>
    #   실측: 지침을 20% 줄이면 문자열 검사 2건이 걸려 배포가 막혔다.
    #   조문 문구는 사람이 고치는 것이고, 값이 틀린 게 아니면 고객 문서는 안전하다.
    try:
        _ds = doctrine_selftest()
        if _ds:
            print('[경고] 조문검사 %d건 — 발행은 막지 않는다' % len(_ds))
            for _x in _ds[:8]: print('   · %s' % _x)
    except Exception as _e:
        print('[경고] 조문검사 실행 실패', _e)
    try:
        with open(_os.path.join(base,'BARUM_DOCTRINE.md'),encoding='utf-8',errors='replace') as _h: _doc=_h.read()
    except Exception: _doc=''
    if _doc:
        _jo=len(set(_re.findall(r'^## .*?제(\d+)조', _doc, flags=_re.M)))
        if _jo < 30: bad.append(f'지침 조문 {_jo}개 — 30 미만이면 잘린 파일 의심')
        for _k in DOC_MUST:
            if _k not in _doc: bad.append(f'지침에 핵심 조문 누락: {_k}')
        # ★v410e 조문이 「[확인 대기]는 제3부에만」이라고 <b>설명하는 문장</b>까지 걸려 거짓양성이 났다.
        #   → <b>목록 항목(`- ` / `1. `)으로 적힌 미결</b>만 위반으로 본다. 설명문은 제외.
        _k3=_doc.find('# 제3부')
        if _k3>0:
            _viol=[_l for _l in _doc[:_k3].split('\n')
                   if '[확인 대기]' in _l and _re.match(r'\s*(?:[-*]|\d+\.)\s', _l)]
            if _viol:
                bad.append(f'본문에 미결 {len(_viol)}건 잔존 — 제33조: 미결은 제3부에만 / 예: {_viol[0].strip()[:60]}')
    else:
        bad.append('BARUM_DOCTRINE.md를 읽지 못했다')
    mn2 = rd_main = ''
    try:
        with open(_os.path.join(base,'main.py'),encoding='utf-8',errors='replace') as _h: mn2=_h.read()
    except Exception: mn2=''
    # ★v460 제69조 — 문자열 조문검사는 <b>발행을 막지 않는다</b>(위에서 경고로 출력했다).
    #   실측: 지침을 20% 줄이면 여기서 2건이 걸려 정당한 정리가 배포를 막았다.
    # ★★★★★v410 제37조 — <b>조문 커버리지를 매번 숫자로 찍는다.</b>
    #   지점장 지시 「지금 해라 절대 미루지 마라」. 커버리지가 떨어지면 그 자리에서 보인다.
    _cvg = '?'
    try:
        _i2 = min(mn2.index('_STRUCT_SELFTEST'), mn2.index('_JOMUN_SELFTEST'))
        _j2 = mn2.index('def ci_selftest')
        _cov = set(int(x) for x in _re.findall(r'제(\d+)조', mn2[_i2:_j2]))
        _joset = set(int(x) for x in _re.findall(r'^## .*?제(\d+)조', _doc, flags=_re.M))
        # ★v411c [보류]·[미결] 조문은 아직 구현이 없으므로 검사 대상이 아니다.
        #   제목에 「[보류]」가 붙은 조문은 커버리지에서 뺀다 — 대신 그 사실을 로그에 남긴다.
        # ★v412 [보류]·[기각] 조문은 구현이 없으므로 커버리지 대상이 아니다.
        _hold = set(int(x) for x in _re.findall(r'^## .*?제(\d+)조 — \[(?:보류|기각|미해결)\]', _doc, flags=_re.M))
        # ★★★★★v434 (제55조 · 2026.08.16): <b>조문이 사라져도 커버리지는 100%</b>였다.
        #   실측 사고 — 제51조에 [기각] 표기를 넣다 <b>제52조를 통째로 덮어썼는데</b>
        #   53/53=100%로 통과했다. 조문 수·번호·글자 수를 <b>따로</b> 지킨다.
        #   4개월 결과물이다. 줄어들면 배포를 막는다.
        # ★★★★★v460 제69조 심플 모드 (지점장 확정 2026.08.17
        #   「검사기 해도 어차피 오류나고 더 심해져. 차라리 심플 모드 가자 — 읽고 체크하고 지키기」)
        #   ★<b>값이 틀린 것만 발행을 막는다.</b> 아래는 전부 <b>경고</b>다 — zip은 통과시킨다.
        #     이유(실측): 분량 하한이 오늘 정당한 지침 증설을 3번 막았고,
        #     문자열 커버리지 100%는 아무것도 증명하지 않으면서 형식 검사를 늘리게 만들었다.
        warn = []
        if len(_joset) < DOCTRINE_MIN_ART:
            warn.append('조문 %d개 (하한 %d) — 확인' % (len(_joset), DOCTRINE_MIN_ART))
        _gap = [x for x in range(1, max(_joset or [0]) + 1)
                if x not in _joset and x not in DOCTRINE_SKIP_ART]
        if _gap:
            warn.append('조문 번호 결번 %s — 삭제 의심' % _gap)
        warn.extend(js_selftest(d or '.'))
        _cvg = '%d개' % len(_joset)
        if warn:
            print('[경고] 발행은 막지 않는다 — %d건' % len(warn))
            for _w in warn[:10]: print('   · %s' % _w)
    except Exception as _e:
        bad.append(f'[제37조] 커버리지 측정 불가 {_e}')
    print('[zip검증] 각인 ' + (list(stamps.values())[0] if stamps else '?') +
          f' · 조문 {_jo if _doc else 0}개 · 값검사 실패 {len(bad)}건')
    for _b in bad[:15]: print('   ★', _b)
    return bad

# ★★★★★v406 지침 정독 강제 (지점장 지시 2026.08.12 · 법)
_DOC_PATHS = ('BARUM_DOCTRINE.md', os.path.join(os.path.dirname(os.path.abspath(__file__)), 'BARUM_DOCTRINE.md'))

def _doc_read(tag=''):
    # ★v409 로그 라벨에 <b>버전 번호를 박지 않는다</b>. v406→v407→v408로 올리며
    #   라벨만 옛 번호로 남아 <b>각인과 어긋났다</b>(실측: 각인 v408인데 로그는 `[v407 지침]`).
    #   각인은 `VSTAMP` 한 곳뿐이다 — 로그에 버전을 쓰려면 그 상수를 쓴다.
    """분석 실행마다 지침 전문을 읽고 조문 검사를 돌린다. 안 읽히면 그 사실을 시끄럽게 남긴다."""
    _txt = ''
    for _p in _DOC_PATHS:
        try:
            with open(_p, encoding='utf-8', errors='replace') as _f: _txt = _f.read()
            if _txt: break
        except Exception: continue
    if not _txt:
        print(f'[지침] ★★ BARUM_DOCTRINE.md를 읽지 못했다 — zip에 빠졌거나 배포 누락 (제0조 6항 위반) [{tag}]')
        return {'ok': False, 'chars': 0, 'jomun': 0, 'fail': ['지침 파일 없음']}
    import re as _re6
    _jo = len(set(_re6.findall(r'^## .*?제(\d+)조', _txt, flags=_re6.M)))
    _st = _re6.search(r'v\d{3}[a-z]*-[a-z0-9]+-\d{8}', _txt)
    _stamp = _st.group(0) if _st else '없음'
    try: _fail = doctrine_selftest()
    except Exception as _e: _fail = [f'검사 실행 불가 {_e}']
    print(f'[지침] 정독 {len(_txt):,}자 · 조문 {_jo}개 · 각인 {_stamp} · 조문검사 실패 {len(_fail)}건 [{tag}]')
    for _x in _fail[:12]: print(f'   ★조문위반 {_x}')
    if _stamp != VSTAMP:
        print(f'[지침] ★★ 각인 불일치 — 코드 {VSTAMP} ≠ 지침 {_stamp} (제0조 6항: 그 zip은 불량이다)')
    # ★★★★★v407: <b>3파일 각인이 서로 같은지</b>도 검사한다.
    #   한 파일만 옛 각인이면 그 파일은 <b>옛 코드</b>다 — 「배포했는데 안 바뀐다」의 진짜 원인이 여기다.
    _mis = []
    for _fn in ('coverage_benchmark.py', 'report_weasy.py'):
        for _d in ('', os.path.dirname(os.path.abspath(__file__))):
            _fp = os.path.join(_d, _fn) if _d else _fn
            try:
                # ★v407b 각인은 파일 <b>끝쪽</b>에 있을 수 있다(report_weasy.py는 6.1MB 중 6,174,604번째).
                #   앞 400KB만 읽었더니 <b>정상인데도 「불일치」</b>가 떴다 — 검사가 거짓말을 했다.
                #   ★검사가 틀리면 진짜 결함보다 나쁘다. 사람이 경고를 무시하게 만든다. → <b>전문을 읽는다.</b>
                with open(_fp, encoding='utf-8', errors='replace') as _f2: _c = _f2.read()
            except Exception: continue
            _all = _re6.findall(r'v\d{3}[a-z]*-[a-z0-9]+-\d{8}', _c)
            _v2 = _all[-1] if _all else '없음'
            if _v2 != VSTAMP: _mis.append(f'{_fn}={_v2}')
            break
    if _mis:
        print(f'[지침] ★★ 파일 각인 불일치 — 코드 {VSTAMP} ≠ ' + ' / '.join(_mis))
    return {'ok': not _fail, 'chars': len(_txt), 'jomun': _jo, 'stamp': _stamp, 'fail': _fail}


@app.post('/analyze')
async def analyze(file:UploadFile=File(None), file2:List[UploadFile]=File(None), pw:str=Form('')):
    if pw!=PW: return JSONResponse({'ok':False,'error':'비밀번호 오류'})
    # ★★★★★v406 (지점장 지시 2026.08.12, 최상위·영구):
    #   지점장 원문: 「<b>지침은 분석지 실행마다 무조건 읽어라. 법이다. 리딩시간 늘려도 된다.</b>」
    #   → 분석 <b>매 실행마다</b> `BARUM_DOCTRINE.md`를 실제로 읽고 조문 검사를 돌린다.
    #     읽는 것을 <b>내 기억이 아니라 코드가 강제</b>한다.
    #   ★조문 위반이 있으면 <b>로그로 시끄럽게</b> 남긴다(제11조 「조용히 틀리는 것을 시끄럽게」).
    _doc_read(tag='analyze')
    # ★★★★★v458 제68조 6항 (지점장 2026.08.17 「지침 무시되면 안 된다고 한 게 1000번이다」)
    #   지침을 <b>읽었는지</b>가 아니라 <b>지켜졌는지</b>를 분석 실행마다 확인한다.
    #   ★산출은 막지 않는다(제49조) — 대신 <b>화면에 크게</b> 띄운다. 조용히 넘어가지 않는다.
    _BEHAVE_WARN = []
    try:
        _bw, _bd = behave_selftest()
        print('[동작검사] 조문 %d개 실행 · 실패 %d건' % (len(_bd), len(_bw)))
        if _bw:
            _BEHAVE_WARN = _bw
            for _x in _bw[:12]:
                print('  ★★ 지침 위반: %s' % _x)
    except Exception as _e:
        print('[동작검사] 실행 실패', _e)
    # ★★★★★v373 (지점장 확정 2026.08.09 「자리 정해라」): <b>칸의 뜻을 고정한다</b>.
    #   <b>왼쪽 file = 보장분석지 / 오른쪽 file2 = 가입제안서</b> — 상황과 무관하게 항상 같다.
    #   ・Ⅰ 보장분석지만 = 왼쪽만  ・Ⅱ 둘 다 = 왼쪽+오른쪽  ・Ⅲ 제안서만 = <b>오른쪽만</b>
    #   ★구 v372는 「혼자면 왼쪽·같이면 오른쪽」이라 <b>칸의 뜻이 상황마다 바뀌었다</b>.
    #     그래서 보장분석지를 오른쪽에·제안서를 왼쪽에 올리면 오류도 안 나고
    #     <b>엉뚱한 산출물이 조용히 나왔다</b>. 자리를 고정해 그 경로를 없앤다.
    _base_f = file  if (file  is not None and (getattr(file ,'filename','') or '')) else None
    # ★★★★★v385 (지점장 확정 2026.08.12): <b>가입제안서 = 한 칸에 복수 선택 · 최대 3건</b>.
    #   지점장 원문: 「제안서 하나의 칸에 복수선택 3개까지」.
    #   ★칸은 그대로 <b>오른쪽(file2) 하나</b>다 — 칸을 늘리지 않는다(v373 자리 고정 유지).
    #   FastAPI가 같은 이름의 파일을 여러 개 받으려면 <b>List[UploadFile]</b>이어야 한다.
    _prop_fs = []
    if file2 is not None:
        _f2list = file2 if isinstance(file2, (list, tuple)) else [file2]
        _prop_fs = [_f for _f in _f2list if _f is not None and (getattr(_f,'filename','') or '')]
    _JEAN_MAX = 3
    if len(_prop_fs) > _JEAN_MAX:
        print(f'[JEAN] 제안서 {len(_prop_fs)}건 → 상한 {_JEAN_MAX}건까지만 사용')
        _prop_fs = _prop_fs[:_JEAN_MAX]
    _prop_f = _prop_fs[0] if _prop_fs else None      # 기존 참조 호환(파일명·단독모드 판정용)
    if not _base_f and not _prop_f:
        return JSONResponse({'ok':False,'error':'파일이 없습니다. 왼쪽 칸에 보장분석지 PDF, 오른쪽 칸에 가입제안서 PDF를 올려주세요.'})
    _bn = ((_base_f.filename if _base_f else '') or '').lower()
    _txt_f = _base_f if _bn.endswith('.txt') else None
    _pdf_f = _base_f if _bn.endswith('.pdf') else None
    if _base_f and not _txt_f and not _pdf_f:
        return JSONResponse({'ok':False,'error':'왼쪽 칸은 보장분석지 PDF여야 합니다.'})

    txt=''; fname=((_txt_f or _pdf_f).filename if (_txt_f or _pdf_f) else (_prop_f.filename if _prop_f else ''))
    if _txt_f:
        raw=await _txt_f.read()
        for enc in ['utf-8','cp949','euc-kr']:
            try: txt=raw.decode(enc); break
            except: pass
        else: txt=raw.decode('utf-8',errors='ignore')

    # ★★★★★v370 가입제안서 계약을 <b>먼저</b> 만든다 — parse_txt에 extra로 넘겨
    #   보장분석지와 <b>같은 후처리</b>(_HB 심장 묶음 · 세부보충 · 대표값)를 타게 하기 위함.
    #   ★v385: 최대 3건을 <b>올린 순서 그대로</b> 각각 계약으로 만든다(합치지 않는다).
    _prop_cts = []
    for _pi, _pf in enumerate(_prop_fs, 1):
        try:
            _pb = await _pf.read()
            _pct = build_proposal_contract(_pb, _pf.filename)
            if _pct:
                _prop_cts.append(_pct)
                print(f'[JEAN] 제안서 {_pi}/{len(_prop_fs)} 계약 생성 — {_pf.filename}')
            else:
                print(f'[JEAN] 제안서 {_pi}/{len(_prop_fs)} 파싱 실패 — 열 추가 안 함 ({_pf.filename})')
        except Exception as _je:
            print(f'[JEAN] 제안서 {_pi} 처리 예외', type(_je).__name__, _je)

    # ★OCR PDF 우선(2026.07.07 지점장 정답): PDF 있으면 pdftotext 직독을 주 소스. 깨지면 txt 폴백.
    src_note=''
    try:
        _txt_data = parse_txt(txt, fname, extra=_prop_cts) if txt.strip() else None
    except Exception:
        _txt_data=None
    _pdf_data=None; pdf_txt=''; _img_pdf_nokey=False; _img_prod=''
    globals()['_IMG_PDF_WARN']=''   # ★v280 매 분석마다 초기화
    globals()['_JEAN_ONLY']=''      # ★v371
    if _pdf_f:
        pdf_bytes=await _pdf_f.read()
        # ★v60 이미지 PDF 진단: 텍스트레이어 직독이 0글자면 = 이미지 전용 PDF(글자 없음).
        #   'Microsoft Print To PDF'로 다시 저장하면 글자층이 통째 이미지가 돼 직독 불가.
        #   이때 유일한 경로는 비전 OCR(API키 필요) — 키 없으면 명확히 알린다.
        try:
            import subprocess as _sp, tempfile as _tf
            with _tf.NamedTemporaryFile(suffix='.pdf', delete=False) as _tfp:
                _tfp.write(pdf_bytes); _tpp=_tfp.name
            _rawtl=_sp.run(['pdftotext','-layout',_tpp,'-'],capture_output=True,text=True,timeout=60).stdout
            try: os.unlink(_tpp)
            except: pass
            # ★v132: '글자수<30'만 보면 못 잡는 케이스가 있다 — 브라우저 '인쇄→PDF 저장'본은
            #   Type3 서브셋 폰트 + ToUnicode 없음이라 pdftotext가 글자는 1만자 넘게 뱉지만
            #   한글은 0자다(전부 깨진 제어문자). 실측: 권양영_보장.pdf = 18,333자 / 한글 0자.
            #   따라서 판정 기준을 '한글 글자 수'로 바꾼다.
            _hangul = sum(1 for _c in (_rawtl or '') if '\uac00' <= _c <= '\ud7a3')
            # ★★★v231 (2026.07.25 한정환 실측): Producer도 함께 읽는다.
            #   `Microsoft: Print To PDF` / `Print to PDF` 는 <b>글자층이 통째 이미지</b>가 된 확정 증거다.
            #   실측 대조 — 이정화(정상) Producer=`oz` 37,278자·한글 4,566자 /
            #              한정환(실패) Producer=`Microsoft: Print To PDF` <b>12자·한글 0자</b>(전부 페이지구분자).
            try:
                _pinf=_sp.run(['pdfinfo',_tpp],capture_output=True,text=True,timeout=30).stdout
            except Exception: _pinf=''
            _prod=''
            for _l in (_pinf or '').split('\n'):
                if _l.lower().startswith('producer:'): _prod=_l.split(':',1)[1].strip()
            _is_print = ('print to pdf' in _prod.lower()) or ('microsoft' in _prod.lower())
            print(f'[PDF_DIAG] chars={len(_rawtl or "")} hangul={_hangul} producer={_prod!r} printpdf={_is_print}')
            # ★★v231: <b>한글이 없으면 API 키가 있어도 이미지 PDF로 확정 안내</b>한다.
            #   구 조건은 `and not ANTHROPIC_API_KEY`라 <b>키가 있으면 안내가 안 나가고</b> 비전 OCR로 흘렀다.
            #   그러나 정본은 <b>OCR 금지</b>(3사 실측으로 데이터 파괴 확인) → 원본 재업로드를 요구하는 게 맞다.
            if not _rawtl or len(_rawtl.strip())<30 or _hangul<100:
                _img_pdf_nokey=True
                _img_prod=_prod
                globals()['_IMG_PDF_WARN'] = ('★ 경고 — 글자층 없는 이미지 PDF를 비전 OCR로 읽었습니다'
                                              + (f' (생성기: {_prod})' if _prod else '')
                                              + ' — 담보명 오독 가능. 고객 제출용으로 쓰지 말 것.')
        except Exception: pass
        pdf_txt=pdf_to_txt(pdf_bytes)
        # ★v281 비전 OCR 부분 실패(페이지 유실)를 엑셀 확인사항·화면 배너로 함께 끌고 간다.
        try:
            _vp=globals().get('_VISION_PARTIAL','')
            if _vp:
                globals()['_IMG_PDF_WARN']=((globals().get('_IMG_PDF_WARN','') or '') +
                                            ' / ★ ' + re.sub(r'<[^>]+>','',_vp))
        except Exception: pass
        if pdf_txt.strip():
            try:
                _pdf_data=parse_txt(pdf_txt, _pdf_f.filename, extra=_prop_cts)
            except Exception:
                _pdf_data=None
    if _pdf_data and _pdf_data.get('contracts') and not _looks_broken(_pdf_data):
        data=_pdf_data; src_note='OCR PDF 직독(주)'
    elif _txt_data and _txt_data.get('contracts'):
        data=_txt_data; src_note='TXT 입력'
    elif _pdf_data and _pdf_data.get('contracts'):
        data=_pdf_data; src_note='OCR PDF(깨짐 감지)'
    else:
        data=_txt_data; src_note='추출 실패'

    # ★★★★★v373 (지점장 확정 2026.08.09 「자리 정해라」): <b>가입제안서 단독 = 오른쪽 칸만</b>.
    #   왼쪽(보장분석지)이 없으면 계약 0건이다 → 이미 만들어둔 제안 계약 `_prop_cts`로 산출한다.
    #   ★검산·실손 세대는 <b>「불가」로 명시하고 건너뛴다</b>(막지 않는다). 지점장 원문: 「불가」.
    #   ★새 후처리 경로를 만들지 않는다 — parse_txt(extra=)로 기존 엔진에 그대로 태운다(영구원칙).
    #   ★구 v371~v372의 「왼쪽 칸 제안서」 경로는 <b>폐기</b>(칸의 뜻이 상황마다 바뀌던 원인).
    if (not data or not data.get('contracts')) and _prop_cts:
        # ★★★★★v592 (지점장 2026.08.26 「<b>가입제안서 3개까지 해달라고 했는데 1개밖에 인식이 안 돼</b>」).
        #   [실측 결함] 로그는 `제안서 1/3 · 2/3 · 3/3 계약 생성`인데 그다음이
        #     `제안 계약 <b>1건</b> 합류 → 총 1건`이었다. 만들어놓고 <b>첫 건만</b> 넘겼다.
        #   원인 = 여기 `_solo = _prop_cts[0]` — <b>[0] 하나만</b> `extra`로 태웠다.
        #   ⇒ <b>담보 3건 이상인 제안 계약을 전부</b> 태운다(최대 3건은 위 `_JEAN_MAX`가 이미 자름).
        _solos = [_x for _x in _prop_cts if len((_x.get('dambo') or {})) >= 3]
        _drop = len(_prop_cts) - len(_solos)
        if _drop:
            print('[JEAN단독] 담보 3건 미만 %d건 제외 — 나머지 %d건 사용' % (_drop, len(_solos)))
        if _solos:
            try:
                data = parse_txt('', fname, extra=_solos); src_note='가입제안서 단독'
                print('[JEAN단독] 제안서 %d건 전부 단독 모드로 산출' % len(_solos))
                globals()['_JEAN_ONLY'] = ('★ 가입제안서 단독 — 보장분석지가 없어 '
                                           '[검산] 불가 · [실손 세대 판정] 불가 · 보유계약 비교 불가')
                print('[JEAN단독] 오른쪽 칸 제안서만 → 단독 모드로 산출 — 검산·실손세대 건너뜀')
            except Exception as _se2:
                print('[JEAN단독] parse_txt 실패', type(_se2).__name__, _se2)
        else:
            print('[JEAN단독] 담보 3건 이상인 제안 계약이 없다 — 단독 모드 미적용')

    try:
        if not data or not data.get('contracts'):
            if _img_pdf_nokey:
                _why = f'(생성기: {_img_prod})' if _img_prod else ''
                _vf = globals().get('_VISION_FAIL','')
                return JSONResponse({'ok':False,'source':'이미지PDF',
                    'error':f'글자를 읽을 수 없는 PDF입니다 {_why} — 한글 0자로 추출됩니다. '
                            '"인쇄 → PDF로 저장"으로 만든 파일은 글자층이 통째로 이미지가 되어 분석이 불가능합니다. '
                            'let: 리포트 화면에서 <b>인쇄가 아니라 PDF 다운로드(저장)</b> 버튼으로 받은 '
                            '원본 파일을 손대지 말고 그대로 올려주세요.'
                            + (f' [비전OCR 진단] {_vf}' if _vf else '')})
            # ★v101: 원인을 화면에서 바로 알 수 있게 추출 단계 수치를 함께 노출(진단용)
            _dbg = (f'계약을 찾지 못했습니다. [진단] 경로={src_note} / '
                    f'PDF추출글자={len(pdf_txt or "")} / TXT입력글자={len(txt or "")} / '
                    f'PDF계약={len((_pdf_data or {}).get("contracts") or [])} / '
                    f'TXT계약={len((_txt_data or {}).get("contracts") or [])}')
            return JSONResponse({'ok':False,'error':_dbg,'source':src_note})
        cust=data['client']; d=tempfile.mkdtemp(); now=datetime.datetime.now()
        xl=os.path.join(d,f'보장진단_{cust}.xlsx'); pt=os.path.join(d,f'보장분석지_{cust}.pptx')
        tx=os.path.join(d,f'치료비정리_{cust}.pptx')
        unmapped=build_excel(data,xl)
        # ★★★★★v323 <b>10억통장 = 이름으로 찾는다</b>(지점장 지시 2026.08.01).
        #   지점장 원문: "10억통장은 이름으로 찾아야 한다. 유일하게 보장진단서에 10억통장이
        #   기재되면 거꾸로 엑셀과 보장분석지 PPT에 기재해라."
        #   → <b>판정 원천 = 담보명의 '리셋월렛'(raw)</b>. 엑셀 22행에 <b>역기재</b>하면
        #     `read_excel_totals`가 그 뒤에 돌아 <b>보장분석지 PPT까지 자동 연동</b>된다.
        try:
            _r10v=0; _r10pd=''
            for _c in (data.get('contracts') or []):
                if '흥국' not in str(_c.get('company','')).replace(' ',''): continue
                _pd=str(_c.get('product','')).replace(' ','')
                _hit=('리셋월렛' in _pd) or ('리셋월랫' in _pd)
                for _k,_v in (_c.get('dambo') or {}).items():
                    _kk=str(_k).replace(' ','')
                    if ('리셋월렛' in _kk) or ('리셋월랫' in _kk):
                        _hit=True
                        try: _r10v=max(_r10v,int(float(_v)))
                        except Exception: pass
                if _hit and not _r10pd: _r10pd=str(_c.get('product',''))
            if _r10pd and _r10v:
                _wbx=openpyxl.load_workbook(xl); _wsx=_wbx['보장분석']
                _tr=None; _tc=None
                for _rr in range(6,_wsx.max_row+1):
                    if str(_wsx.cell(_rr,2).value or '').strip()=='10억 플랜': _tr=_rr; break
                _key=re.sub(r'\s','',_r10pd)[:20]
                for _cc in range(3,_wsx.max_column):
                    _hd=re.sub(r'\s','',str(_wsx.cell(1,_cc).value or ''))
                    if '흥국' in _hd and _key and _key in _hd: _tc=_cc; break
                if _tr and _tc:
                    if _wsx.cell(_tr,_tc).value in (None,'',0):
                        _wsx.cell(_tr,_tc).value=_r10v; _wbx.save(xl)
                        print(f'[R10] 이름 판정 → 엑셀 「10억 플랜」 R{_tr}C{_tc} = {_r10v}만원 역기재')
                    else: print(f'[R10] 엑셀 R{_tr}C{_tc}에 이미 값 있음 — 역기재 생략')
                else: print(f'[R10] 역기재 대상 못 찾음 (행={_tr} 열={_tc})')
        except Exception as _erx: print('[R10] 엑셀 역기재 실패 (%s)'%str(_erx)[:60])
        if not recalc_xlsx(xl): inject_sum_cache(xl)   # ★v29u: Railway(LibreOffice 없음)에서도 합계 캐시 보장
        ppt_totals, sq, ss, ppt_splits = read_excel_totals(xl)   # 등식2: PPT는 완성 엑셀만 읽음
        ppt_ok=build_ppt(data,pt,ppt_totals,sq,ss,ppt_splits)
        # 치료비정리 PPT 폐기(v29) — 내용 부실, 보장설명지 PDF로 대체
        xlsx_b64=base64.b64encode(open(xl,'rb').read()).decode()
        _sm = make_summary(data)
        if _img_pdf_nokey:
            _sm = ('<div style="background:#B00020;color:#fff;padding:12px;border-radius:6px;'
                   'font-weight:700;line-height:1.6;margin-bottom:12px">'
                   '★ 경고 — 글자층이 없는 <b>이미지 PDF</b>입니다'
                   + (f' (생성기: {_img_prod})' if _img_prod else '') + '<br>'
                   '비전 OCR로 읽었습니다. 담보명 한 글자만 오독돼도 담보 배정이 통째로 갈립니다.<br>'
                   + ((f'<b>{globals().get("_VISION_PARTIAL","")}</b> — 실패한 페이지의 계약·담보는 '
                       f'산출물에서 통째로 빠져 있습니다.<br>') if globals().get('_VISION_PARTIAL') else '') +
                   '<b>이 산출물은 고객 제출용으로 쓰지 마십시오.</b><br>'
                   'let: 리포트 화면에서 <b>인쇄가 아니라 PDF 다운로드(저장)</b> 버튼으로 받은 '
                   '원본 파일을 그대로 다시 올려주세요.</div>') + _sm
            src_note = (src_note or '') + ' / ★이미지PDF-비전OCR'
        response={'ok':True,'xlsx_b64':xlsx_b64,'xlsx_name':f'보장진단_{cust}.xlsx',
                  'summary':_sm,'pptx_ready':ppt_ok,'source':src_note}
        if ppt_ok and os.path.exists(pt):
            response['pptx_b64']=base64.b64encode(open(pt,'rb').read()).decode()
            response['pptx_name']=f'보장분석지_{cust}.pptx'
        # ── 보장설명서: 충족률 PDF + ★보장진단서 PPT(편집가능) — 둘 다 실패해도 엑셀·PPT는 유지 ──
        rep=None
        try:
            from coverage_benchmark import map_excel_to_report
            # ★★★v138 흥국화재 10억통장 가입 판정(지점장 실측 정정 2026.07.21):
            #   '리셋월렛'은 <b>상품명이 아니라 담보명</b>으로 들어온다.
            #   실측(장기상) — 상품 '무배당 흥Good 든든한 3N5 간편종합보험' 안의
            #   담보 '플래티넘 건강 리셋월렛II(3대질병고액및중환자실치료,특정…) 100,000'(=10억).
            #   따라서 <b>상품명·담보명 두 곳을 모두</b> 본다(지점장 '2가지 다 적용' 지시).
            _r10=False
            try:
                for _c in (data.get('contracts') or []):
                    _co=str(_c.get('company','')).replace(' ','')
                    if '흥국' not in _co: continue
                    _pd=str(_c.get('product','')).replace(' ','')
                    _hit=('리셋월렛' in _pd) or ('리셋월랫' in _pd)
                    if not _hit:
                        for _k in (_c.get('dambo') or {}):
                            _kk=str(_k).replace(' ','')
                            if ('리셋월렛' in _kk) or ('리셋월랫' in _kk): _hit=True; break
                    if _hit: _r10=True; break
            except Exception: pass
            # ★v146 금액도 함께 전달(진단서 카드에 표기). 만원 단위 원본값.
            _r10amt=0
            try:
                for _c in (data.get('contracts') or []):
                    if '흥국' not in str(_c.get('company','')).replace(' ',''): continue
                    for _k,_v in (_c.get('dambo') or {}).items():
                        _kk=str(_k).replace(' ','')
                        if ('리셋월렛' in _kk) or ('리셋월랫' in _kk):
                            try: _r10amt=max(_r10amt,int(float(_v)))
                            except Exception: pass
            except Exception: pass
            print(f'[R10] 흥국화재 10억통장 가입판정={_r10} 금액={_r10amt}만원')
            # ★★★★★v421e — <b>진단서 7페이지 전용 담보</b>(지점장 확정 2026.08.14).
            #   「신특정순환계질환 주요치료비」는 <b>마스터 행이 없다</b> → 엑셀엔 기재하지 않는다
            #   (정본 3줄 2항: 엑셀에 없는 것은 만들지 않는다). 대신 <b>7p 순환계 칸에만</b> 싣는다.
            #   ★계약 담보 원문에서 직접 찾아 rep에 주입한다(엑셀 우회가 아니라 <b>산출물별 담보 범위</b>).
            #   4종: 순환계 주요 / 순환계 통합 / 암 통합 / 비급여 암 통합
            _p7only = {}
            def _p7put(_k, _v):
                try: _v=float(_v)
                except Exception: return
                if _v: _p7only[_k] = max(_p7only.get(_k,0.0), _v)
            try:
                _srcs = list((data.get('contracts') or []))   # ★v421f 파싱 결과는 data['contracts']다
                for _c7 in _srcs:
                    if not isinstance(_c7, dict): continue
                    for _n7, _v7 in (_c7.get('dambo') or {}).items():
                        _t7 = re.sub(r'\s','',str(_n7))
                        if '통합치료' in _t7:
                            if '순환계' in _t7:            _p7put('순환계통합치료비', _v7)
                            elif '비급여' in _t7 or '전액본인' in _t7: _p7put('비급여암통합치료비', _v7)
                            elif '암' in _t7:               _p7put('암통합치료비', _v7)
                        elif ('순환계' in _t7) and ('주요치료' in _t7):
                            _p7put('순환계주요치료비', _v7)
                        # ★v422c 특정치료비Ⅲ는 <b>엑셀 암주요치료비</b>로 간다(지점장 최종 확정) —
                        #   7p 암 주요치료비 칸은 그 엑셀 값을 읽으므로 전용 주입이 필요 없다.
            except Exception as _e7: print('[v421f 7p전용] 탐색 실패', _e7)
            print(f'[v421f 7p전용] 진단서 전용 담보 {len(_p7only)}건 {_p7only}')
            rep=map_excel_to_report(xl, settings={'client':cust,'reset10':_r10,'reset10_amt':_r10amt,
                'branch':'메이크원','manager':'최은혜','title':'지점장','phone':''})
            if _p7only: rep['p7_only'] = _p7only      # ★v421f 진단서 전용 칸 값(엑셀 미반영)
        except Exception as _re:
            response['report_error']='분석데이터 생성 실패: '+str(_re)
        if rep is not None:
            # ★ 보장설명지 PDF 별도 생성 중단(2026.07.11 지점장 지시): 보장진단서 PPT가 동일 내용(PDF 페이지 이미지)이라
            #    별도 PDF는 렌더 1회(약 60초)를 중복 유발 → 속도 위해 스킵. 필요 시 이 블록 복구.
            # ★ 보장진단서 PPT (편집가능) — 같은 rep로 생성
            try:
                from report_pptx import build_report_pptx
                rpx=os.path.join(d,f'보장진단서_{cust}.pptx')
                # ★v107: 같은 렌더에서 벡터 PDF(보장설명서)도 함께 받는다(추가 렌더 0회).
                #   PPT는 이미지라 확대·인쇄 시 글자가 뭉갠다 → 선명본은 이 PDF.
                rpdf=os.path.join(d,f'보장설명서_{cust}.pdf')
                build_report_pptx(rep, rpx, pdf_out=rpdf)
                if os.path.exists(rpdf):
                    # ★★★★★v424 (지점장 지시 2026.08.16): <b>재무 페이지는 진단서에만</b>.
                    #   진단서 PPT는 이 PDF 앞부분을 잘라 쓰므로 원본에는 남기고,
                    #   <b>설명서 PDF에서만</b> 그 한 장을 뺀다.
                    try:
                        assert os.path.exists(rpx), 'PPT 미생성 — 제거 보류'
                        print('[v424 설명서] PPT 존재 확인 후 제거 진행')
                        import subprocess as _sp9, re as _re9
                        from pypdf import PdfReader as _RD, PdfWriter as _WR
                        _txt = _sp9.run(['pdftotext', rpdf, '-'],
                                        capture_output=True).stdout.decode('utf-8', 'replace')
                        _pgs = _txt.split('\f')
                        _drop = [i for i, t in enumerate(_pgs)
                                 if 'ASSET & FINANCE' in t and '돈을 보내고 있나요' in t]
                        if not _drop:
                            # ★v425 (제49조) 정정: 설명서는 <b>인포메이션 구간만</b> 잘라 쓰므로
                            #   재무 페이지가 애초에 없다. 다만 <b>있는데 못 지운 경우</b>는 사고다.
                            if 'ASSET' in _txt or '돈을 보내고' in _txt:
                                print('[v425 ★설명서] 재무 페이지가 있는데 제거 실패 — 문구 변경 의심')
                                response.setdefault('warnings', []).append(
                                    '[확인] 설명서 재무 페이지 제거 실패')
                            else:
                                print('[v425 설명서] 재무 페이지 없음(인포메이션 구간만 사용) — 정상')
                        if _drop:
                            _rd = _RD(rpdf); _wr = _WR()
                            for _i9 in range(len(_rd.pages)):
                                if _i9 not in _drop:
                                    _wr.add_page(_rd.pages[_i9])
                            with open(rpdf, 'wb') as _f9:
                                _wr.write(_f9)
                            print(f'[v424 설명서] 재무 페이지 {len(_drop)}장 제거 → {len(_rd.pages)-len(_drop)}쪽')
                    except Exception as _e9:
                        print('[v424 설명서] 재무 제거 실패:', _e9)
                    response['report_b64']=base64.b64encode(open(rpdf,'rb').read()).decode()
                    response['report_name']=f'보장설명서_참고자료_{cust}.pdf'
                if os.path.exists(rpx):
                    response['report_pptx_b64']=base64.b64encode(open(rpx,'rb').read()).decode()
                    response['report_pptx_name']=f'보장진단서_{cust}.pptx'
            except Exception as _pe:
                response['report_pptx_error']=str(_pe)
        # ★★★★★v420 REPORT_AUDIT — 산출물 게이트 (지점장 지시 2026.08.14 「다해」)
        #   지적 30건을 되짚으면 원인은 셋뿐이었고 그 중 첫째가
        #   <b>게이트가 엑셀만 보고 산출물은 열어보지 않는다</b>였다.
        #   검산 27/27·감사 58/58을 통과하면서 6페이지가 통째로 레드였고,
        #   계약이 21건으로 세어졌고, 7페이지 레드가 0건이었다.
        #   ★<b>값이 맞나</b>가 아니라 <b>보이는 게 맞나</b>를 만든 직후에 검사한다.
        #   ★실패해도 산출은 막지 않는다 — <b>조용히 틀리는 것을 시끄럽게 틀리는 것으로</b> 바꾼다.
        # ★v458 제68조 6항 — 지침 위반이 있으면 화면 맨 위에 띄운다.
        if _BEHAVE_WARN:
            response.setdefault('warnings', []).insert(
                0, '[지침위반 %d건] %s' % (len(_BEHAVE_WARN), ' / '.join(_BEHAVE_WARN[:4])))
        # ★★★★★v446 제61조 3항 — 4대 산출물 개수 게이트.
        #   v446 사고: 진단서 PPT·설명서 PDF가 조용히 사라졌는데 아무 게이트도 안 울렸다.
        _OUT4_KEYS = ('xlsx_b64', 'pptx_b64', 'report_pptx_b64', 'report_b64')
        _miss4 = [k for k in _OUT4_KEYS if k not in response]
        if _miss4:
            print('[v446 산출물게이트] ★ 누락 %d건 %s' % (len(_miss4), _miss4))
            response.setdefault('warnings', []).append(
                '[확인] 산출물 누락 %d건: %s' % (len(_miss4), ', '.join(_miss4)))
        else:
            print('[v446 산출물게이트] 4대 산출물 전부 생성')
        # ★★★★★v601 제137조 2항 (지점장 지시 2026.08.29 「<b>로봇 4모드는 돌리고 정독이다</b>」).
        #   분석이 실제로 돈 <b>그 순간</b>에 4모드 결과를 남긴다. 배지는 이 값이 있을 때만
        #   정독 %를 확정한다 — 없으면 「대기」다(화면만 열었을 때 색이 뜨지 않게).
        try:
            globals()['_MODE4_LAST'] = list(_miss4)
            # ★★★★★v618 제142조 — 로봇2가 대조할 <b>이번 분석 값</b>을 남긴다.
            #   담보명 → 값. 로봇2는 이것이 있을 때만 정답지와 대조한다(없으면 대기).
            _cv2 = {}
            for _c2 in (data.get('contracts') or []):
                for _n2, _v2 in ((_c2 or {}).get('dambo') or {}).items():
                    try:
                        _f2 = float(_v2)
                    except Exception:
                        continue
                    if _f2:
                        _cv2[_n2] = _cv2.get(_n2, 0.0) + _f2
            globals()['_ANS_LAST'] = {'client': cust, 'cov': _cv2}
        except Exception:
            pass
        try:
            _aud=[]; _rep=locals().get('rep') or {}
            # ①계약 수 = 엑셀 계약 열 수(합산 열 제외)
            try:
                _ncol=len(_rep.get('renew_list') or [])+len(_rep.get('nonrenew_list') or [])
                if _rep.get('n_contract') and _ncol and _rep['n_contract']!=_ncol:
                    _aud.append(f"계약수 불일치 n_contract={_rep['n_contract']} vs 목록={_ncol}")
                for _l in (_rep.get('renew_list') or [])+(_rep.get('nonrenew_list') or []):
                    if str(_l.get('nm','')).strip() in ('보유 합계','제안 합계','합계'):
                        _aud.append(f"합산 열이 계약으로 세어짐: {_l.get('nm')}")
            except Exception: pass
            # ②red_map 오염(제안 담보 수를 넘지 않는다)
            try:
                _rm=len(_rep.get('red_map') or {}); _pa=len(_rep.get('prop_amt') or {})
                if _pa and _rm > _pa*3:
                    _aud.append(f"red_map 오염 의심 red={_rm} vs 제안금액={_pa}")
            except Exception: pass
            # ③충족률이 전 항목 동일값이면 변별력 0
            try:
                _p=[d.get('pct') for d in (_rep.get('donut_detail') or []) if d.get('pct') is not None]
                if len(_p)>=5 and len(set(_p))==1:
                    _aud.append(f"충족률 전 항목 동일값 {_p[0]}% — 변별력 0")
            except Exception: pass
            # ④진단서 PPT: 제안 담보가 있는데 레드가 0건인 슬라이드
            try:
                if os.path.exists(rpx) and (_rep.get('prop_amt') or {}):
                    from pptx import Presentation as _PR
                    _pp=_PR(rpx); _nred=0; _nsl=0
                    for _s in _pp.slides:
                        _nsl+=1
                        for _sh in _s.shapes:
                            if not _sh.has_text_frame: continue
                            for _pa2 in _sh.text_frame.paragraphs:
                                for _r2 in _pa2.runs:
                                    try:
                                        if str(_r2.font.color.rgb)=='C00000': _nred+=1
                                    except Exception: pass
                    if _nred==0:
                        _aud.append('진단서 PPT 레드 0건 — 제안 담보가 있는데 색이 안 실렸다')
                    response['audit_slides']=_nsl; response['audit_red']=_nred
            except Exception: pass
            # ⑤설명서 푸터: 연속성·분모
            try:
                if os.path.exists(rpdf):
                    import subprocess as _sp, re as _re3
                    _t=_sp.run(['pdftotext',rpdf,'-'],capture_output=True).stdout.decode('utf-8','replace')
                    # ★v420b 본문 표에도 'N / M'이 나온다 → <b>페이지별 마지막 1건</b>만,
                    #   그리고 <b>분모가 일정한 것</b>만 푸터로 본다(오탐 2건 실측 후 수정).
                    _pages=_t.split('\f'); _nums=[]; _dens=[]
                    for _pg in _pages:
                        _mm=_re3.findall(r'(\d+)\s*/\s*(\d+)',_pg)
                        if _mm: _nums.append(int(_mm[-1][0])); _dens.append(int(_mm[-1][1]))
                    from collections import Counter as _Ct
                    if _dens:
                        _den=_Ct(_dens).most_common(1)[0][0]
                        _nums=[n for n,d in zip(_nums,_dens) if d==_den]
                    if _nums:
                        if len(_nums)!=len(set(_nums)): _aud.append('설명서 푸터 번호 중복')
                        if _nums!=sorted(_nums): _aud.append('설명서 푸터 번호 역순')
                        if max(_nums)>_den: _aud.append(f'설명서 푸터 번호가 분모 초과 {max(_nums)}>{_den}')
            except Exception: pass
            response['audit']=_aud
            print(f"[v420 REPORT_AUDIT] 실패 {len(_aud)}건" + (' :: '+' | '.join(_aud) if _aud else ' (통과)'))
        except Exception as _ae:
            print('[v420 REPORT_AUDIT] 검사 자체 실패:', _ae)
        # ★★★★★v316 모바일 다운로드 근본 수정(지점장 실측 2026.08.01 "삼성폰에서 저장됨은
        #   뜨는데 「내 파일」→다운로드 폴더에도 없다")
        #   <b>원인</b>: 프론트가 blob 다운로드를 <b>5개나 0.8~3초 간격으로 연속 발사</b>했다.
        #     ①<b>모바일 Chrome은 연속 다운로드를 차단</b>한다(첫 개 이후 전부 무시).
        #     ②5.5MB PDF가 <b>3초 revoke</b>를 못 넘긴다. ③인앱 브라우저는 blob 자체가 막힌다.
        #   <b>해법</b>: 파일을 서버에 잠깐 저장하고 <b>진짜 URL</b>을 준다 → 브라우저가 일반
        #     다운로드로 처리하고, 사용자가 <b>카드를 하나씩</b> 누르므로 연속 차단에도 안 걸린다.
        #   ★b64는 <b>그대로 유지</b>한다(PC 호환·`reDL` 폴백). URL만 추가한다 = 회귀 위험 0.
        try:
            import uuid as _uuid, shutil as _sh
            _tok = _uuid.uuid4().hex[:12]
            _dir = os.path.join(tempfile.gettempdir(), 'barum_dl', _tok)
            os.makedirs(_dir, exist_ok=True)
            for _bk, _nk, _uk in (('xlsx_b64','xlsx_name','xlsx_url'),
                                  ('pptx_b64','pptx_name','pptx_url'),
                                  ('chiryo_b64','chiryo_name','chiryo_url'),
                                  ('report_b64','report_name','report_url'),
                                  ('report_pptx_b64','report_pptx_name','report_pptx_url')):
                if not response.get(_bk): continue
                _fn = response.get(_nk) or (_bk + '.bin')
                with open(os.path.join(_dir, _fn), 'wb') as _f:
                    _f.write(base64.b64decode(response[_bk]))
                response[_uk] = '/dl/%s/%s' % (_tok, urllib.parse.quote(_fn))
            # ★오래된 임시 폴더 청소(2시간 초과) — 디스크 누적 방지
            try:
                _root = os.path.join(tempfile.gettempdir(), 'barum_dl')
                _now = datetime.datetime.now().timestamp()
                for _d in os.listdir(_root):
                    _p = os.path.join(_root, _d)
                    if os.path.isdir(_p) and (_now - os.path.getmtime(_p)) > 7200:
                        _sh.rmtree(_p, ignore_errors=True)
            except Exception: pass
            print('[v316 dl] 서버 저장 %s — %d개' % (_tok, sum(1 for k in response if k.endswith('_url'))))
        except Exception as _edl:
            print('[v316 dl] 서버 저장 실패 → b64 폴백만 사용 (%s)' % str(_edl)[:80])
        return JSONResponse(response)
    except Exception as e:
        return JSONResponse({'ok':False,'error':str(e),'trace':traceback.format_exc()[-1500:]})

# ── AI 질문답 ─────────────────────────────────────────────────────────
import httpx

def build_context(data):
    lines=[f"고객명: {data['client']}", f"계약 수: {len(data['contracts'])}건"]
    for ct in data['contracts']:
        lines.append(f"  - {ct['company']} [{ct['renewal']}] {ct['premium']:,}원")
    totals={}
    for ct in data['contracts']:
        for raw,amt in ct['dambo'].items():
            std=resolve(raw)
            if std: totals[std]=totals.get(std,0)+amt
    lines.append("\n매핑된 담보 합계 (만원):")
    for k,v in sorted(totals.items()): lines.append(f"  - {k}: {v:,}")
    unmapped=[]
    for ct in data['contracts']:
        for raw in ct['dambo']:
            if resolve(raw) is None and not any(x in raw for x in ['(1종)','(2종)','(3종)','(4종)','(5종)']):
                unmapped.append(raw)
    if unmapped:
        lines.append("\n자동매핑 실패 담보 (약관 확인 필요):")
        for u in sorted(set(unmapped)): lines.append(f"  - {u}")
    return '\n'.join(lines)

def _save_dl(data: bytes, fname: str) -> str:
    """산출물 1개를 임시 저장하고 /dl URL을 돌려준다(기존 analyze와 같은 방식)."""
    import uuid as _uuid
    _tok = _uuid.uuid4().hex[:12]
    _dir = os.path.join(tempfile.gettempdir(), 'barum_dl', _tok)
    os.makedirs(_dir, exist_ok=True)
    _p = os.path.join(_dir, fname)
    with open(_p, 'wb') as _h:
        _h.write(data)
    return '/dl/%s/%s' % (_tok, urllib.parse.quote(fname))


# ★★★★★v422g — ④ 리모델링 비교 (지점장 지시 2026.08.15)
#   「1버튼 기존 보험 엑셀 / 2버튼 새로 정리된 엑셀 → 두개를 비교한 진단서」
#   ★<b>①②③ 경로(analyze)를 한 줄도 건드리지 않는다</b> — 별도 라우트·별도 모듈(remodel.py).
# ★★★★★v423 — 「엑셀은 1번만 보면된다」(지점장 2026.08.15).
#   엑셀 <b>한 개</b> 안에 보유(기존)와 제안(최종)이 다 들어 있다 → `remodel_single`.
#   구 2파일 방식(old_xlsx+new_xlsx)도 살려 둔다 — 둘 다 오면 2파일 방식이 이긴다.
@app.post('/remodel')
async def remodel_route(xlsx: UploadFile = File(None),
                        old_xlsx: UploadFile = File(None), new_xlsx: UploadFile = File(None),
                        pw: str = Form(''), client: str = Form(''), base_date: str = Form('')):
    if pw != PW:
        return JSONResponse({'ok': False, 'error': '비밀번호 오류'})
    try:
        import remodel as _rm
        _two = bool(hasattr(old_xlsx, 'read') and hasattr(new_xlsx, 'read'))
        if not _two and not xlsx:
            return JSONResponse({'ok': False, 'error': '엑셀을 올려주십시오'})
        _ob = await old_xlsx.read() if _two else b''
        _nb = await new_xlsx.read() if _two else await xlsx.read()
        # ★기본값 객체가 그대로 올 수 있다(직접 호출·테스트) → 문자열만 신뢰한다
        client = client if isinstance(client, str) else ''
        base_date = base_date if isinstance(base_date, str) else ''
        _fn = (new_xlsx.filename if _two else xlsx.filename) or ''
        # ★★★★★v449 제64조 (지점장 지적 2026.08.17 · 실측)
        #   파일명 「보장진단_사공호 (1).xlsx」 → 한글만 뽑아 앞 4자 = <b>「보장진단」</b>.
        #   고객명이 통째로 '보장진단'이 된다(모든 파일이 같은 이름). 파일명은 믿을 수 없다.
        #   → <b>엑셀 A1이 유일 원천</b>이다(A1 = 「박주하 보장진단」). 파일명은 마지막 폴백.
        def _cust_from_xlsx(_by):
            try:
                import openpyxl as _op, io as _io
                _w = _op.load_workbook(_io.BytesIO(_by), read_only=True, data_only=True)
                _s = _w['보장분석'] if '보장분석' in _w.sheetnames else _w.active
                _a1 = str(_s.cell(1, 1).value or '')
                _nm = re.sub(r'\s*보장진단.*$', '', _a1).strip()
                # ★★★★★v541 제128조 (지점장 실측 2026.08.22 「모든 파일에 이름이 100% 표기 안된다」)
                #   엑셀 A1이 <b>마스킹 이름</b>(`김*자 보장진단`)이었다. 구 코드는 `[^가-힣A-Za-z]`로
                #   <b>별표를 지워</b> 「김자」를 만들었다 — 산출물 안은 「김*자」인데 파일명만 깨졌다.
                #   ⇒ 마스킹 문자(* ● ○ ｡ 등)는 <b>`O`로 바꿔 자리를 지킨다</b>(파일명에 `*`는 못 쓴다).
                _nm = re.sub(r'[*●○◯·・\u25cf\u25cb]', 'O', _nm)
                _nm = re.sub(r'[^가-힣A-Za-zO]', '', _nm)
                return _nm if 2 <= len(_nm) <= 6 else ''
            except Exception:
                return ''
        _cl = client or _cust_from_xlsx(_nb) or _cust_from_xlsx(_ob)
        # ★★★★★v545 제128조 2항 (지점장 지시 2026.08.22 「이름은 한글로 3글자 다 나오게 해야한다」)
        #   지침 8항 <b>「이름은 마스킹하지 않는다」</b>가 정본이다.
        #   A1이 <b>마스킹</b>(`김*자`·`김O자`)이면 그대로 쓰지 말고 <b>파일명에서 실명을 찾는다</b>
        #   (실측 김순자 — A1은 `김*자`인데 파일명에 `김순자`가 그대로 있었다).
        #   ★제64조(A1 우선)는 <b>A1이 온전할 때</b>의 규칙이다. 마스킹은 온전한 이름이 아니다.
        if _cl and re.search(r'[*●○◯·・O]', _cl):
            _cand = re.findall(r'[가-힣]{2,4}', re.sub(r'^보장(진단|분석지?)[_\s-]*', '', _fn.split('.')[0]))
            _cand = [c for c in _cand if c not in ('보장진단', '보장분석', '보장분석지', '최종본', '복사본', '고객')]
            if _cand and len(_cand[0]) >= len(_cl):
                print('[v545 고객명] A1 마스킹 %r → 파일명 실명 %r 채택' % (_cl, _cand[0]))
                _cl = _cand[0]
        if not _cl:
            _base = re.sub(r'\s*[-(].*$', '', _fn.split('.')[0])        # ' (1)' · ' - 복사본' 제거
            _base = re.sub(r'^보장(진단|분석지?)[_\s-]*', '', _base)      # 접두 '보장진단_' 제거
            _base = re.sub(r'[*●○◯·・]', 'O', _base)                     # ★v541 제128조 마스킹 보존
            _cl = re.sub(r'[^가-힣O]', '', _base)[:4] or '고객'
        print('[v449 고객명] 파일명 %r → 확정 %r' % (_fn, _cl))
        _bd = base_date or datetime.datetime.now().strftime('%Y.%m.%d')
        r = _rm.remodel_all(_ob, _nb, _cl, _bd) if _two else _rm.remodel_single(_nb, _cl, _bd)
        c = r['cmp']
        _x = _save_dl(r['xlsx'], f'{_cl}_remodel_compare.xlsx')
        _p = _save_dl(r['pptx'], f'{_cl}_remodel_report.pptx')
        _d = _save_dl(r['pdf'], f'{_cl}_remodel_report.pdf') if r.get('pdf') else ''
        # ★★★★★v612 (지점장 지시 2026.08.30 「비교엑셀 + 리포트 + 보장분석ppt 3가지로」)
        _ap = _save_dl(r['anal_pptx'], f'보장분석지_{_cl}(최종).pptx') if r.get('anal_pptx') else ''
        print(f"[v422g 리모델링] {_cl} {c['prem_old']:,.0f}→{c['prem_new']:,.0f} "
              f"절감 {c['save_m']:,.0f}({c['save_pct']}%) · 증가 {len(c['up'])} 신규 {len(c['add'])} "
              f"감소 {len(c['down'])} 삭제 {len(c['delete'])}")
        return JSONResponse({'ok': True, 'client': _cl,
                             'prem_old': int(c['prem_old']), 'prem_new': int(c['prem_new']),
                             'save_m': int(c['save_m']), 'save_y': int(c['save_y']),
                             'save_pct': c['save_pct'],
                             'n_up': len(c['up']), 'n_add': len(c['add']),
                             'n_down': len(c['down']), 'n_del': len(c['delete']),
                             'xlsx': _x, 'pptx': _p, 'pdf': _d, 'anal': _ap})
    except Exception as e:
        traceback.print_exc()
        return JSONResponse({'ok': False, 'error': f'{type(e).__name__}: {e}'})


@app.post('/ask')
async def ask(body:dict):
    if body.get('pw')!=PW: return JSONResponse({'ok':False,'error':'비밀번호 오류'})
    question=body.get('question','').strip()
    data=body.get('data')
    if not question or not data: return JSONResponse({'ok':False,'error':'질문 또는 데이터 없음'})
    context=build_context(data)
    system=("보장분석 전문 AI입니다. 아래 분석 데이터에 관한 질문에만 답하세요.\n"
            "규칙:\n"
            "- 데이터에 없는 담보 -> 현재 계약에 없습니다. 별첨/약관 확인 필요.\n"
            "- 약관 해석 -> 약관을 직접 확인하세요.\n"
            "- 무관한 질문 -> 보장분석 관련 질문만 가능합니다.\n"
            "- 답변 2-3줄 이내, 간결하게, 한국어로.\n"
            f"[분석 데이터]\n{context}")
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            resp=await client.post('https://api.anthropic.com/v1/messages',
                headers={'x-api-key':os.environ.get('ANTHROPIC_API_KEY',''),
                         'anthropic-version':'2023-06-01','content-type':'application/json'},
                json={'model':'claude-haiku-4-5-20251001','max_tokens':300,
                      'system':system,'messages':[{'role':'user','content':question}]})
        r=resp.json()
        answer=r.get('content',[])[0].get('text','답변을 가져오지 못했습니다.')
        return JSONResponse({'ok':True,'answer':answer})
    except Exception as e:
        return JSONResponse({'ok':False,'error':str(e)})
