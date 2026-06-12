# -*- coding: utf-8 -*-
"""
fill_ppt_v2 — 빈 마스터(MASTER_blank.pptx)에 값만 끼워넣는다.
규칙(드리프트 차단):
  · tf.clear() 절대 금지. 기존 run의 .text만 교체 → 폰트·크기·색·줄간격 보존.
  · 라벨 재타이핑 금지. 라벨은 빈 마스터에 이미 있다. 콜론 뒤 빈 run에만 값 주입.
  · 모르면 빈칸(법칙18 추측금지). 매핑 안 된 칸은 손 안 댐.
입력: let_engine.한장보장표(d) 의 53담보 dict.
"""
import sys, re, fitz
from pptx import Presentation
import let_engine as E

MASTER = 'MASTER_보장분석지_PPT_빈폼.pptx'

def won(v):
    v=int(v or 0)
    if v<=0: return ''
    eok,man=divmod(v,10000)
    if eok and man: return f'{eok}억{man:,}만'
    if eok: return f'{eok}억'
    return f'{man:,}만'

def _do(p, full, label, value):
    colon=full.find(':', full.find(label)+len(label))
    if colon<0: return False
    run_of=[]
    for ri,r in enumerate(p.runs): run_of += [ri]*len(r.text)
    cr=run_of[colon] if colon<len(run_of) else len(p.runs)-1
    for ri in range(cr+1, len(p.runs)):
        p.runs[ri].text=value; return True
    p.runs[cr].text=p.runs[cr].text+value
    return True

def inject(shape, label, value):
    """label 뒤 첫 콜론 다음 빈 run에 value 주입. 1패스 시작앵커(반깁스 차단)→2패스 어디든(다필드 일당)."""
    if not value: return False
    ps=[p for p in shape.text_frame.paragraphs if p.runs]
    for p in ps:                                   # 1패스: 문단 시작 앵커
        full=''.join(r.text for r in p.runs)
        if full.lstrip().startswith(label) and _do(p,full,label,value): return True
    for p in ps:                                   # 2패스: 어디든(중간 라벨)
        full=''.join(r.text for r in p.runs)
        if label in full and _do(p,full,label,value): return True
    return False

def setbox(shape, value):
    """라벨 없는 값박스(뇌26 등): 첫 run에 값만."""
    if not value: return False
    tf=shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text=value; return True
    return False

# (박스, 라벨) -> 엔진키  : 확실히 단일라벨인 칸만. 애매한 칸(운전자·일당·치매 등)은 비움.
MAP = [
 ('TextBox 11','상해사망','상해사망'),
 ('TextBox 10','종신','질병사망'),
 ('TextBox 8','질병 3%','질병후유3%'),('TextBox 8','질병 80%','질병후유80%'),
 ('TextBox 8','상해3%','상해후유3%'),('TextBox 8','상해 80%','상해후유80%'),
 ('TextBox 14','암진단비','일반암'),('TextBox 14','전이암','전이암'),
 ('TextBox 14','고액암','고액암'),('TextBox 14','항암치료','고액항암'),
 ('TextBox 49','산정특례','중증뇌혈관'),
 ('TextBox 56','산정특례','중증심장'),
 ('TextBox 17','질병수술','질병수술'),('TextBox 17','뇌혈관수술','뇌혈관수술'),
 ('TextBox 17','심장 수술','심장수술'),
 ('TextBox 19','상해수술','상해수술'),
 ('TextBox 7','골절','골절진단'),('TextBox 7','화상','화상진단'),
 ('TextBox 7','깁스','깁스치료'),('TextBox 7','응급실','응급실내원'),
 ('TextBox 5','일상배상책임','일상생활배상'),
 ('TextBox 6','입원','질병입원의료'),('TextBox 6','통원','질병통원의료'),
 ('TextBox 22','질병일당','질병입원'),('TextBox 22','상해일당','상해입원'),
 ('TextBox 22','질병중환자실','질병중환자실'),('TextBox 22','상해중환자실','상해중환자실'),
]
BOXMAP = [('TextBox 26','뇌혈관진단')]  # 뇌 단일 값박스(검증됨)

def fill(slide, name, date, t):
    by={sh.name:sh for sh in slide.shapes if sh.has_text_frame}
    if 'TextBox 1' in by: setbox(by['TextBox 1'], name)
    for b,key in [('TextBox 16',date[0]),('TextBox 24',date[1]),('TextBox 25',date[2])]:
        if b in by and key: setbox(by[b], key)
    done=0
    for box,label,key in MAP:
        if box in by and inject(by[box], label, won(t.get(key,0))): done+=1
    for box,key in BOXMAP:
        if box in by and setbox(by[box], won(t.get(key,0))): done+=1
    # 심 진단비(법칙71): 허혈성=허혈성칸(54) / 급성심근=급성심근칸(55). 배너박스31 미사용.
    from pptx.util import Pt
    hy=won(t.get('허혈성심장',0)); gs=won(t.get('급성심근경색',0))
    for n,lab,val in [('TextBox 54','허혈성',hy),('TextBox 55','급성심근',gs)]:
        if n in by and val:
            tf=by[n].text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                r=tf.paragraphs[0].runs[0]; r.text=f'{lab} {val}'; r.font.size=Pt(9)
            tf.word_wrap=False; done+=1
    # 한 칸짜리 박스 줄바꿈 방지(겹침 차단)
    for nm in ['TextBox 1','TextBox 16','TextBox 24','TextBox 25','TextBox 26']:
        if nm in by:
            try: by[nm].text_frame.word_wrap=False
            except Exception: pass
    return done

def customer(pdf):
    d=fitz.open(pdf); t,_=E.한장보장표(d)
    p0=d[0].get_text('text')
    name='고객'
    for ln in [x.strip() for x in p0.split('\n') if x.strip()]:
        m=re.match(r'(.+?)\s*고객님', ln)
        if m: name=m.group(1).strip(); break
    md=re.search(r'(\d{4})\.(\d{2})\.(\d{2})', p0)
    date=(md.group(1)[2:],md.group(2),md.group(3)) if md else ('','','')
    return name,date,t

def run(pdf):
    name,date,t=customer(pdf)
    prs=Presentation(MASTER)
    n=fill(prs.slides[0], name, date, t)
    out=f'/mnt/user-data/outputs/보장분석지_{name}.pptx'
    prs.save(out); print('saved',out,'| 채운칸',n)
    return out

if __name__=='__main__': run(sys.argv[1])
