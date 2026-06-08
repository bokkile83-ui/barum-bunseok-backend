# -*- coding: utf-8 -*-
import engine_v2 as E
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

BLUE=RGBColor(0,0,0xFF); BLACK=RGBColor(0,0,0); RED=RGBColor(0xFF,0,0)
cell=E.cell; COLS=E.COLS; 종수술=E.종수술; 종수술g=E.종수술g; STAT=E.STAT
_116d=getattr(E,'_116d',{})

def sz_of(para):
    for rr in para.runs:
        if rr.font.size: return rr.font.size
    return None
def red(para):
    for rr in para.runs:
        try: rr.font.color.rgb=RED
        except: pass
def put(para, rows):                      # 비갱신(검정)/갱신(파랑)
    if isinstance(rows,int): rows=[rows]
    b=sum(cell[r][c][0] for r in rows for c in COLS)
    g=sum(cell[r][c][1] for r in rows for c in COLS)
    if b==0 and g==0: red(para); return
    sz=sz_of(para)
    if b:
        run=para.add_run(); run.text=f' {b}'; run.font.color.rgb=BLACK
        if sz: run.font.size=sz
    if g:
        run=para.add_run(); run.text=(f'/{g}' if b else f' {g}'); run.font.color.rgb=BLUE
        if sz: run.font.size=sz
def jong(para, r, n):
    tot=[0]*6; anyg=False; hit=False
    for (rk,cc),d in 종수술.items():
        if rk==r:
            hit=True
            for k in range(1,7): tot[k-1]+=d.get(str(k),0)
            anyg|=종수술g[(rk,cc)]
    if not hit: red(para); return
    para.runs[0].text='( '+'/'.join(str(x) for x in tot[:n])+' )'
    for rr in para.runs[1:]: rr.text=''
    for rr in para.runs: rr.font.color.rgb=(BLUE if anyg else BLACK)
def n대(para):
    tot=[0]*6; anyg=False
    for col,d in _116d.items():
        for k,v in d.items(): tot[k-1]+=v
        if STAT[col][0]=='갱신': anyg=True
    if not any(tot): red(para); return
    sz=sz_of(para)
    run=para.add_run(); run.text=' '+'/'.join(str(x) for x in tot[:6])
    run.font.color.rgb=(BLUE if anyg else BLACK)
    if sz: run.font.size=sz

M={'뇌혈관':35,'뇌졸증':36,'뇌출혈':37,'허혈성':49,'급성심근':50,
 '산정특례(뇌혈관)':41,'산정특례(심장)':47,
 '유사암':25,'항암치료':34,'표적치료':26,'암통원비':33,'다빈치로봇수술비':32,
 '질병수술':72,'뇌혈관수술':76,'허혈성수술':77,'심장 수술':79,'5대기관수술':82,
 '상해수술':64,'골절수술':69,'5대골절수술':70,'화상수술':71,'중대상해수술':65,'창상봉합수술':68,
 '골절':84,'5대골절':85,'중대화상':89,'반깁스':90,'깁스':91,'응급실':86,
 '대인':92,'대물':93,'합의금':94,'6주미만':95,'변호사비':96,'자부상':97,
 '입원':98,'통원':99,'MRI':101,
 '질병 3%':18,'질병 80%':19,'상해3%':16,'상해 80%':17,
 '질병사망':12,'종신':10,'상해사망':15,
 '질병일당':53,'상해일당':54,'질병중환자실':62,'상해중환자실':63,'암일당':33,
 '간병인일당':56,'간호통합병동일당':61,'임플란트':87,
 '1인실 상급병원일당':57,'1인실 종합병원일당':58}
BOX={('TextBox 49','혈전용해치료비'):42,('TextBox 49','2대주요치료비'):40,
 ('TextBox 56','혈전용해치료비'):52,('TextBox 56','2대주요치료비'):48}

p=Presentation('폼_보장분석지_17p.pptx')
x=p.slides._sldIdLst; sl=list(x)
for i in range(len(sl)-1,0,-1): x.remove(sl[i])

def fill(shapes):
    for sh in shapes:
        if sh.shape_type==6: fill(sh.shapes); continue
        if not sh.has_text_frame: continue
        sh.text_frame.word_wrap=True
        try: sh.text_frame.auto_size=MSO_AUTO_SIZE.SHRINK_TEXT_ON_OVERFLOW
        except: pass
        bn=sh.name
        for pi,para in enumerate(sh.text_frame.paragraphs):
            if not para.runs: continue
            txt=''.join(r.text for r in para.runs); key=txt.split(':')[0].strip()
            if txt.strip().startswith('(') and '/' in txt:
                row=74 if bn=='TextBox 17' else 67 if bn=='TextBox 19' else None
                if row: jong(para,row,5)
                continue
            if '대 수술' in txt and bn=='TextBox 17' and pi==4: n대(para); continue
            if key=='암진단비': put(para,[22,23,20,21]); continue
            if key=='심장': put(para,[43,44,45,46]); continue
            if '양성자' in txt and '표적' in txt:
                bg=lambda r:(sum(cell[r][c][0] for c in COLS),sum(cell[r][c][1] for c in COLS))
                t=sum(bg(26)); pp=sum(bg(27)); ss=sum(bg(28))
                para.runs[0].text=f"표적: {t or ''} / 양성자·세기: {pp or ''}/{ss or ''}"
                for rr in para.runs[1:]: rr.text=''
                continue
            row=BOX.get((bn,key))
            if row is None: row=M.get(key)
            if row is None: continue
            put(para,row)

fill(p.slides[0].shapes)
# 칸 넘침 방지: 폰트 직접 축소 (제목/날짜 제외)
from pptx.util import Pt
_skip={'TextBox 21','TextBox 29','TextBox 35','TextBox 36'}
def shrink(shapes):
    for sh in shapes:
        if sh.shape_type==6: shrink(sh.shapes); continue
        if not sh.has_text_frame or sh.name in _skip: continue
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                if r.font.size and r.font.size.pt>8:
                    r.font.size=Pt(max(8,round(r.font.size.pt*0.8)))
shrink(p.slides[0].shapes)
for sh in p.slides[0].shapes:
    if sh.shape_type==6:
        for s2 in sh.shapes:
            if s2.has_text_frame and '님의 보장' in s2.text_frame.text and s2.text_frame.paragraphs[0].runs:
                s2.text_frame.paragraphs[0].runs[0].text='두*호'+s2.text_frame.paragraphs[0].runs[0].text
                s2.text_frame.word_wrap=False   # 제목 한 줄
p.save('보장분석지_두호.pptx')
print('done')
