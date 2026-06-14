# -*- coding: utf-8 -*-
import os, io, json, base64, zipfile, re, tempfile, traceback, datetime
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import JSONResponse, HTMLResponse
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

CATS=[
('사망',['일반사망(종신)','질병사망(80세)','질병사망(__세)','상해사망','교통상해사망']),
('후유장애',['상해후유3%','상해후유80%','질병후유3%','질병후유80%','교통후유장해']),
('암',['암진단비','암진단비(갱신)','유사암(갑.기.경.제)','3대암진단비','통합전이암진단비','중입자치료비','항암치료비','표적항암치료비','다빈치로봇수술비','양성자','세기조절','암수술','암주요치료비(상급병원)','하이클래스(비급여주요치료비)','암통원비(상급병원)','암일당','CAR-T세포치료비','13대항암치료비','면역항암치료비','호르몬항암치료비']),
('뇌혈관',['뇌혈관진단비','뇌졸증진단비','뇌졸증진단비(갱신)','뇌출혈진단비','산정특례(뇌혈관)','혈전용해치료비(뇌졸중)','순환계주요치료비','뇌혈관수술비']),
('심장',['부정맥질환(I49)','심부전','염증','협심증','빈맥','심근병증','심장판막','급성심근경색','허혈성','산정특례(심장)','순환계주요치료비','심장수술비']),
('일당',['질병일당','상해일당','질병중환자실','상해중환자실','1인실 상급병원일당','1인실 종합병원일당','간병인사용일당','간병인지원일당(질병)','간병인지원일당(상해)','요양병원일당','간호통합병동일당','암일당(입원)']),
('입원비',['상해입원','질병입원']),
('수술비',['상해수술비','질병수술비','7대질병 수술비','상해종수술비(1-5종)','질병종수술비(1-5종)','119대 수술비','골절수술비','5대골절수술비','중대상해수술비','화상수술비','창상봉합수술비']),
('운전자',['교통사고처리지원금','6주미만처리지원금','교통사고벌금(대인)','교통사고벌금(대물)','변호사선임비용','자동차부상위로금','일상생활배상책임']),
('골절',['골절진단비','골절진단비(치아제외)','5대골절진단비','5대골절수술비']),
('화상',['화상진단비','중대화상진단비','응급실내원비(응급)','깁스치료비']),
('화재/치아',['화재벌금/화재보험','크라운','임플란트','치아 기타']),
('기타',['독감','간병인지원일당','재가급여','시설급여']),
('실손',['실손입원','실손통원','실손 약(처방조제)','MRI','도수치료','비급여주사'])]

BLUE='FF0000FF';RED='FFFF0000';WHITE='FFFFFFFF';BLACK='FF000000';GREEN='FF006400';GRAY='FFEAEAEA'
thin=Side('thin',color='FF999999');med=Side('medium',color='FF333333')
GRID=Border(top=med,bottom=med,left=thin,right=thin)

# ── 제외 판정 ──────────────────────────────────────────────────────────
EXCLUDE=['실효','미납해지','농업인NH안전보험','자동차보험']
def is_excluded(company,product):
    return any(kw in company+product for kw in EXCLUDE)

# ── 갱신 판정 ──────────────────────────────────────────────────────────
def judge_renewal(product,expiry,pay_count):
    if '갱신형' in product and '비갱신' not in product: return '갱신'
    if '갱신' in product and '비갱신' not in product: return '갱신'
    if expiry.startswith('9999'): return '비갱신(종신)'
    try:
        a,b=pay_count.split('/')
        if int(b.strip())>240: return '비갱신'
    except: pass
    return '비갱신'

def get_종번호(name):
    for i,k in enumerate(['(1종)','(2종)','(3종)','(4종)','(5종)'],1):
        if k in name: return i
    return 0

# ── TXT 파싱 ──────────────────────────────────────────────────────────
def parse_txt(txt):
    lines=[l.rstrip() for l in txt.replace('\r\n','\n').replace('\r','\n').split('\n')]
    client='고객'
    for l in lines[:30]:
        l=l.strip()
        m=re.match(r'^([가-힣]{2,5})\s*$',l)
        if m and len(m.group(1))<=4: client=m.group(1); break
        m2=re.search(r'([가-힣]{2,4})\s+고객님',l)
        if m2: client=m2.group(1); break

    contracts=[]; i=0; n=len(lines)
    while i<n:
        l=lines[i].strip()
        if '실효계약 리스트' in l or '미납해지' in l: break
        if '정상계약 리스트' not in l: i+=1; continue
        i+=1
        while i<n and not lines[i].strip(): i+=1
        if i>=n: break
        company=lines[i].strip(); i+=1
        if is_excluded(company,''):
            while i<n and '정상계약 리스트' not in lines[i] and '실효계약 리스트' not in lines[i]: i+=1
            continue
        contract_date=expiry_date=pay_period=pay_count=''; premium=0
        for j in range(i,min(i+5,n)):
            l=lines[j]
            m=re.search(r'(\d{4}\.\d{2}\.\d{2})\s*~\s*(\d{4}\.\d{2}\.\d{2})',l)
            if m: contract_date=m.group(1); expiry_date=m.group(2)
            m2=re.search(r'(\d+)\s*/\s*(\d+)\s*회',l)
            if m2: pay_count=f"{m2.group(1)}/{m2.group(2)}"
            m3=re.search(r'([\d,]+)원',l)
            if m3:
                v=int(m3.group(1).replace(',',''))
                if 1000<v<5000000: premium=v
            m4=re.search(r'(\d+)년납',l)
            if m4: pay_period=f"{m4.group(1)}년납"
            m5=re.search(r'월납\s*/\s*(\d+)년',l)
            if m5 and not pay_period: pay_period=f"{m5.group(1)}년납"
        while i<n and not lines[i].strip(): i+=1
        product=''
        for j in range(i,min(i+6,n)):
            l=lines[j].strip()
            if l and not re.search(r'계약자|납입주기|보험료|보장기간',l):
                if len(l)>5 and not re.search(r'^\d+$',l) and not re.search(r'^\d{4}\.\d{2}',l):
                    product=l; i=j+1; break
        renewal=judge_renewal(product,expiry_date,pay_count)
        dambo={}; j=i
        while j<n:
            l=lines[j].strip()
            if '정상계약 리스트' in l or '실효계약 리스트' in l: i=j; break
            m=re.search(r'^(.+?)\s{2,}([\d,]+)\s*$',l)
            if not m: m=re.search(r'^(.+?)\s+([\d,]+)\s*$',l)
            if m:
                name=re.sub(r'\s+',' ',m.group(1).strip())
                try:
                    amt=int(m.group(2).replace(',',''))
                    if 0<amt<=100000 and len(name)>2: dambo[name]=dambo.get(name,0)+amt
                except: pass
            j+=1
        else: i=j
        if company:
            contracts.append({'company':company,'product':product,'contract_date':contract_date,
                'expiry_date':expiry_date,'premium':premium,'pay_period':pay_period,
                'pay_count':pay_count,'renewal':renewal,'dambo':dambo})
    return {'client':client,'contracts':contracts}

# ── 담보명 표준화 ──────────────────────────────────────────────────────
DMAP={
    '상해사망(갱신형) [보통약관]':'상해사망','상해사망':'상해사망','일반상해사망':'상해사망',
    '기본계약(상해사망(간편가입Ⅲ))담보':'상해사망','질병사망':'질병사망(80세)',
    '상해후유장해3%':'상해후유3%','상해후유80%':'상해후유80%','질병후유장해3%':'질병후유3%','질병후유80%':'질병후유80%',
    '일반암진단비':'암진단비','암진단비':'암진단비','암진단Ⅱ(유사암제외)(간편가입Ⅲ)담보':'암진단비',
    '고액암진단비':'3대암진단비',
    '갑상선암.기타피부암.유사암진단비Ⅲ':'유사암(갑.기.경.제)','유사암진단비':'유사암(갑.기.경.제)',
    '유사암진단Ⅱ(양성뇌종양포함)(간편가입Ⅲ)담보':'유사암(갑.기.경.제)',
    '표적항암약물허가치료비':'표적항암치료비','표적항암약물허가치료(간편가입Ⅲ)(갱신형)담보':'표적항암치료비',
    '항암방사선.약물치료비':'항암치료비','항암방사선치료(간편가입Ⅲ)담보':'항암치료비','항암약물치료(간편가입Ⅲ)담보':'항암치료비',
    '항암방사선(세기조절)치료(간편가입Ⅲ)(갱신형)담보':'세기조절','항암세기조절방사선치료비':'세기조절',
    '항암방사선(양성자)치료(간편가입Ⅲ)(갱신형)담보':'양성자','항암양성자방사선치료비':'양성자',
    '암수술(간편가입Ⅲ)담보':'암수술','카티(CAR-T)항암약물허가치료비':'CAR-T세포치료비',
    '뇌혈관질환진단비Ⅲ(건강맞춤형Ⅱ)(갱신형)':'뇌혈관진단비','뇌혈관질환진단비':'뇌혈관진단비',
    '뇌혈관질환진단(간편가입Ⅲ)담보':'뇌혈관진단비',
    '뇌졸중진단비':'뇌졸증진단비','뇌졸중진단(간편가입Ⅲ)담보':'뇌졸증진단비','뇌졸중진단비(건강맞춤형Ⅱ)(갱신형)':'뇌졸증진단비',
    '뇌출혈진단':'뇌출혈진단비',
    '중증질환자(뇌혈관질환)산정특례대상진단비(연간1회한)(건강맞춤형Ⅱ)(갱신형)':'산정특례(뇌혈관)',
    '뇌혈관질환수술비Ⅲ(건강맞춤형Ⅱ)(갱신형)':'뇌혈관수술비','심뇌혈관질환수술(간편가입Ⅲ)담보':'뇌혈관수술비',
    '뇌경색증(I63)혈전용해치료비':'혈전용해치료비(뇌졸중)','혈전용해치료비Ⅱ(뇌졸중)(간편가입Ⅲ)담보':'혈전용해치료비(뇌졸중)',
    '허혈심장질환진단비Ⅲ(건강맞춤형Ⅱ)(갱신형)':'허혈성','허혈심장질환진단비':'허혈성','허혈심장질환진단(간편가입Ⅲ)담보':'허혈성',
    '급성심근경색증진단':'급성심근경색','급성심근경색증진단(간편가입Ⅲ)담보':'급성심근경색',
    '중증질환자(심장질환)산정특례대상진단비(연간1회한)(건강맞춤형Ⅱ)(갱신형)':'산정특례(심장)',
    '허혈심장질환수술비Ⅲ(건강맞춤형Ⅱ)(갱신형)':'심장수술비','허혈심장질환수술비':'심장수술비',
    '급성심근경색증(I21)혈전용해치료비':'혈전용해치료비(뇌졸중)',
    '질병수술비':'질병수술비','질병수술(간편가입Ⅲ)담보':'질병수술비',
    '상해수술비(건강맞춤형Ⅱ)(갱신형)':'상해수술비','상해수술(간편가입Ⅲ)담보':'상해수술비',
    '골절수술(간편가입Ⅲ)담보':'골절수술비','화상수술(간편가입Ⅲ)담보':'화상수술비',
    '120대질병수술Ⅱ(간편가입Ⅲ)(질병수술3(24대질병))담보':'119대 수술비',
    '5대기관질병수술(관혈/비관혈)(연간1회한)(간편가입Ⅲ)담보':'심장수술비',
    '중대한특정상해수술(간편가입Ⅲ)담보':'중대상해수술비',
    '질병입원의료비':'실손입원','상해입원의료비':'실손입원','질병외래의료비':'실손통원',
    '도수/체외충격파/증식치료':'도수치료','비급여주사제':'비급여주사','MRI검사의료비':'MRI',
    '간병인사용질병입원일당(1일이상)(요양병원)(간편가입)(갱신형)':'간병인사용일당',
    '간호간병통합서비스질병입원일당(1-180일)(간편가입)(갱신형)':'간호통합병동일당',
    '상급종합병원질병입원일당(상급병실(1인실),1일이상60일한도)(간편가입)(갱신형)':'1인실 상급병원일당',
    '종합병원질병입원일당(상급병실(1인실),1일이상30일한도)(간편가입)(갱신형)':'1인실 종합병원일당',
    '교통사고처리지원금':'교통사고처리지원금','교통사고벌금(대물)':'교통사고벌금(대물)',
    '교통사고벌금(대인)':'교통사고벌금(대인)','변호사선임비용':'변호사선임비용',
    '무보험차에 의한 상해':'일상생활배상책임',
    '골절진단(간편가입Ⅲ)담보':'골절진단비(치아제외)','골절진단비':'골절진단비(치아제외)',
    '깁스치료담보':'깁스치료비','깁스치료':'깁스치료비',
    '가족생활배상책임':'일상생활배상책임','일상생활배상책임':'일상생활배상책임',
    '보험료납입지원':None,
    '치과치료(보존치료)':'크라운','치과치료(보철치료)':'임플란트',
}
def resolve(raw):
    if raw in DMAP: return DMAP[raw]
    for k,v in DMAP.items():
        if k in raw: return v
    if any(k in raw for k in ['(1종)','(2종)','(3종)','(4종)','(5종)']): return None
    return raw

# ── 엑셀 생성 ──────────────────────────────────────────────────────────
def build_excel(data, path):
    contracts=data['contracts']
    n_ct=len(contracts)
    # 컬럼 번호: 3번(C)부터 계약 순서대로 연속 배치 (빈 열 없음)
    col_indices={i: 3+i for i in range(n_ct)}  # 계약idx → 열번호
    # ★ 합계열 = 마지막 계약 열 + 1 (컴팩트 배치이므로 항상 맨 끝)
    last_col = 3 + n_ct

    # 담보 매트릭스 구성
    matrix={}
    surg15={'질병':{},'상해':{}}

    for idx,ct in enumerate(contracts):
        col_key=idx  # 정수 키로 통일
        dambo=ct['dambo']

        q종={k:v for k,v in dambo.items() if '질병수술비' in k and get_종번호(k)>0}
        s종={k:v for k,v in dambo.items() if '상해수술비' in k and get_종번호(k)>0}
        if q종:
            vals=[0]*5
            for k,v in q종.items():
                i=get_종번호(k)-1
                if 0<=i<5: vals[i]=v
            surg15['질병'][col_key]=vals
        if s종:
            vals=[0]*5
            for k,v in s종.items():
                i=get_종번호(k)-1
                if 0<=i<5: vals[i]=v
            surg15['상해'][col_key]=vals

        skip=set(list(q종.keys())+list(s종.keys()))
        for raw,amt in dambo.items():
            if raw in skip: continue
            std=resolve(raw)
            if std is None: continue
            cat_found=None
            for cat,rows in CATS:
                if std in rows: cat_found=cat; break
            if not cat_found: cat_found='기타'
            key=(cat_found,std)
            if key not in matrix: matrix[key]={}
            matrix[key][col_key]=matrix[key].get(col_key,0)+amt

    # 엑셀 작성
    wb=Workbook(); ws=wb.active; ws.title='보장진단'
    cust=data['client']
    ws['A1']=f"{cust} 고객님\n보장진단"
    ws.merge_cells('A1:B7')
    ws['A1'].alignment=Alignment(horizontal='center',vertical='center',wrap_text=True)
    ws['A1'].font=Font(bold=True,size=12)

    total_premium=0
    for idx,ct in enumerate(contracts):
        cn=col_indices[idx]
        gen=ct['renewal']=='갱신'
        paid='완납' in ct['renewal']
        fill=GREEN if paid else (BLUE if gen else RED)
        h=ws.cell(row=1,column=cn,value=f"{ct['company']}\n[{ct['renewal']}]")
        h.fill=PatternFill('solid',fgColor=fill)
        h.font=Font(color=WHITE,bold=True,size=9)
        h.alignment=Alignment(horizontal='center',vertical='center',wrap_text=True)
        h.border=GRID
        meta=[ct['premium'],ct['contract_date'],ct['expiry_date'],ct['pay_period'],ct['pay_count']]
        labels=['보험료','가입일','만기일','납입기간','납입횟수']
        for i,v in enumerate(meta):
            cc=ws.cell(row=2+i,column=cn,value=v)
            cc.alignment=Alignment(horizontal='center'); cc.border=GRID; cc.font=Font(size=8)
        total_premium+=ct['premium']

    # ★ 합계열 — 마지막 열 (last_col = 3 + n_ct, 항상 데이터 열 뒤)
    hc=ws.cell(row=1,column=last_col,value='합계')
    hc.fill=PatternFill('solid',fgColor='FF2E75B6')
    hc.font=Font(color=WHITE,bold=True,size=9)
    hc.alignment=Alignment(horizontal='center')
    hc.border=GRID
    tp=ws.cell(row=2,column=last_col,value=total_premium)
    tp.font=Font(bold=True,size=8)
    tp.number_format='#,##0'
    tp.border=GRID

    # 담보 행
    # SUM 수식용 열 문자 계산
    def col_letter(n):
        s=''
        while n>0:
            n,r=divmod(n-1,26)
            s=chr(65+r)+s
        return s

    first_data_col=col_letter(3)           # C
    last_data_col=col_letter(last_col-1)   # 마지막 계약열
    sum_col=col_letter(last_col)           # 합계열

    r=8
    for cat,rows in CATS:
        extra=sorted({k[1] for k in matrix if k[0]==cat and k[1] not in rows})
        allrows=rows+extra; start=r
        for rw in allrows:
            ws.cell(row=r,column=2,value=rw).font=Font(size=9)
            ws.cell(row=r,column=2).border=GRID
            ws.cell(row=r,column=1).border=GRID
            for idx in range(n_ct):
                ws.cell(row=r,column=col_indices[idx]).border=GRID
            q=ws.cell(row=r,column=last_col)
            q.border=GRID; q.font=Font(bold=True,size=8)

            if rw in ('상해종수술비(1-5종)','질병종수술비(1-5종)'):
                kind='상해' if '상해' in rw else '질병'
                tot=[0]*5; kk=surg15.get(kind,{})
                for idx in range(n_ct):
                    seq=kk.get(idx)
                    if seq:
                        for j in range(5): tot[j]+=seq[j]
                        cell=ws.cell(row=r,column=col_indices[idx],value='/'.join(str(x) for x in seq))
                        cell.font=Font(color=(BLUE if contracts[idx]['renewal']=='갱신' else BLACK),size=8)
                        cell.alignment=Alignment(horizontal='center',shrink_to_fit=True)
                if any(tot): q.value='/'.join(str(x) for x in tot); q.number_format='General'
            else:
                has=False
                for idx in range(n_ct):
                    amt=matrix.get((cat,rw),{}).get(idx)
                    if amt:
                        has=True
                        cell=ws.cell(row=r,column=col_indices[idx],value=amt)
                        cell.number_format='#,##0'
                        cell.font=Font(color=(BLUE if contracts[idx]['renewal']=='갱신' else BLACK),size=8)
                        cell.alignment=Alignment(horizontal='center')
                # ★ SUM 수식: 법칙 22조 — 하드코딩 금지, 수식 유지
                if has: q.value=f'=SUM({first_data_col}{r}:{last_data_col}{r})'
            r+=1
        ws.merge_cells(start_row=start,start_column=1,end_row=r-1,end_column=1)
        a=ws.cell(row=start,column=1,value=cat)
        a.font=Font(bold=True,size=9)
        a.fill=PatternFill('solid',fgColor=GRAY)
        a.alignment=Alignment(horizontal='center',vertical='center'); a.border=GRID

    ws.column_dimensions['A'].width=6; ws.column_dimensions['B'].width=22
    for idx in range(n_ct):
        ws.column_dimensions[col_letter(col_indices[idx])].width=7.5
    ws.column_dimensions[col_letter(last_col)].width=12

    # 확인사항 시트
    ms=wb.create_sheet('📋확인사항')
    ms['A1']=f"{cust} 고객님 — 확인사항"
    ms['A1'].font=Font(bold=True,size=11)
    ms.append([]); ms.append(['구분','대상','내용'])
    ms.append(['완료',f'계약 {n_ct}건',f'월보험료합계: {total_premium:,}원'])
    ms.column_dimensions['A'].width=12; ms.column_dimensions['B'].width=40; ms.column_dimensions['C'].width=50
    wb.save(path)

# ── PPT 채우기 ────────────────────────────────────────────────────────
HERE=os.path.dirname(os.path.abspath(__file__))
TPL_PPT=os.path.join(HERE,'MASTER_보장분석지_PPT_빈폼.pptx')

def _shrink_box(shape):
    """텍스트 오버플로우 방지: 박스 내 자동 축소"""
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    tf=shape.text_frame
    tf.word_wrap=False
    # autofit=spAutoFit 적용 (XML 직접)
    txBody=tf._txBody
    bodyPr=txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        # 기존 autofit 제거 후 normAutofit(shrink) 추가
        for tag in ['a:noAutofit','a:spAutoFit','a:normAutofit']:
            el=bodyPr.find(qn(tag))
            if el is not None: bodyPr.remove(el)
        norm=etree.SubElement(bodyPr, qn('a:normAutofit'))
        norm.set('fontScale','100000')
        norm.set('lnSpcReduction','0')

def build_ppt(data, path):
    if not os.path.exists(TPL_PPT): return False
    prs=Presentation(TPL_PPT)
    sl=prs.slides[0]
    by={sh.name:sh for sh in sl.shapes if sh.has_text_frame}
    cust=data['client']
    contracts=data['contracts']
    now=datetime.datetime.now()

    # 합계 계산
    totals={}
    surg_q=[0]*5; surg_s=[0]*5
    for ct in contracts:
        for raw,amt in ct['dambo'].items():
            i=get_종번호(raw)-1
            if i>=0:
                if '질병수술비' in raw: surg_q[i]+=amt
                if '상해수술비' in raw: surg_s[i]+=amt
            std=resolve(raw)
            if std: totals[std]=totals.get(std,0)+amt

    def g(nm): return totals.get(nm,0)
    def r_set(box,pi,ri,val):
        if box not in by: return
        tf=by[box].text_frame
        if pi<len(tf.paragraphs):
            p=tf.paragraphs[pi]
            if ri<len(p.runs): p.runs[ri].text=val

    # ★ 오버플로우 방지: 내용이 긴 박스들에 shrink 적용
    overflow_boxes=['TextBox 14','TextBox 17','TextBox 19','TextBox 22','TextBox 57']
    for bname in overflow_boxes:
        if bname in by:
            try: _shrink_box(by[bname])
            except: pass

    # 제목/날짜
    by['TextBox 21'].text_frame.word_wrap=False
    by['TextBox 21'].text_frame.paragraphs[0].runs[0].text=f'{cust} 님의 보장'
    by['TextBox 21'].text_frame.paragraphs[0].runs[1].text='(전)'
    by['TextBox 36'].text_frame.paragraphs[0].runs[0].text=f'{now.year}년'
    by['TextBox 35'].text_frame.paragraphs[0].runs[0].text=f'{now.month:02d}월'
    by['TextBox 29'].text_frame.paragraphs[0].runs[0].text=f'{now.day:02d}일 기준'
    for b in ['TextBox 36','TextBox 35','TextBox 29']:
        by[b].text_frame.word_wrap=False

    # 산정특례 줄바꿈 방지
    for b in ['TextBox 49','TextBox 56']:
        if b in by: by[b].text_frame.word_wrap=False

    # 사망
    q_d=g('질병사망(80세)'); s_d=g('상해사망')
    if q_d: r_set('TextBox 10',2,2,f': {q_d:,}')
    if s_d: r_set('TextBox 11',0,1,f': {s_d:,}')
    종신_d=0
    for ct in contracts:
        if '종신' in ct['renewal']:
            for raw,v in ct['dambo'].items():
                if resolve(raw)=='상해사망': 종신_d+=v
    if 종신_d: r_set('TextBox 10',3,1,f': {종신_d:,}')

    # 후유장해
    if g('상해후유3%'): r_set('TextBox 8',2,1,f'3% : {g("상해후유3%"):,}')
    if g('질병후유3%'): r_set('TextBox 8',0,1,f'3% : {g("질병후유3%"):,}')
    if g('상해후유80%'): r_set('TextBox 8',3,1,f'80% : {g("상해후유80%"):,}')
    if g('질병후유80%'): r_set('TextBox 8',1,1,f'80% : {g("질병후유80%"):,}')

    # 뇌혈관
    if g('뇌혈관진단비'):
        by['TextBox 46'].text_frame.paragraphs[0].runs[0].text=f'뇌혈관\n{g("뇌혈관진단비"):,}'
    if g('뇌졸증진단비'):
        by['TextBox 47'].text_frame.paragraphs[0].runs[0].text=f'뇌졸증\n{g("뇌졸증진단비"):,}'
    if g('뇌출혈진단비'):
        by['TextBox 48'].text_frame.paragraphs[0].runs[0].text=f'뇌출혈\n{g("뇌출혈진단비"):,}'
    if g('산정특례(뇌혈관)'): r_set('TextBox 49',0,3,f': {g("산정특례(뇌혈관)"):,}')
    if g('혈전용해치료비(뇌졸중)'): r_set('TextBox 49',1,1,f': {g("혈전용해치료비(뇌졸중)"):,}')

    # 심장
    if g('허혈성'):
        by['TextBox 54'].text_frame.paragraphs[0].runs[0].text=f'허혈성\n{g("허혈성"):,}'
    if g('급성심근경색'):
        by['TextBox 55'].text_frame.paragraphs[0].runs[0].text=f'급성심근\n{g("급성심근경색"):,}'
    if g('산정특례(심장)'): r_set('TextBox 56',0,3,f': {g("산정특례(심장)"):,}')

    # 암
    if g('암진단비'): r_set('TextBox 14',0,1,f': {g("암진단비"):,}')
    if g('유사암(갑.기.경.제)'): r_set('TextBox 14',1,2,f': {g("유사암(갑.기.경.제)"):,}')
    if g('항암치료비'): r_set('TextBox 14',4,1,f': {g("항암치료비"):,} / ')
    if g('표적항암치료비'): r_set('TextBox 14',5,1,f': {g("표적항암치료비"):,} / ')
    if g('세기조절'): r_set('TextBox 14',5,4,f': {g("세기조절"):,}')
    if g('양성자'): r_set('TextBox 14',5,5,f': {g("양성자"):,}')
    if g('다빈치로봇수술비'): r_set('TextBox 14',7,1,f': {g("다빈치로봇수술비"):,}')

    # 수술비
    if g('질병수술비'): r_set('TextBox 17',0,1,f': {g("질병수술비"):,}')
    if any(surg_q):
        r_set('TextBox 17',3,0,f'({"/".join(str(x) for x in surg_q)})')
        r_set('TextBox 17',3,2,'')
    if g('뇌혈관수술비'): r_set('TextBox 17',5,1,f': {g("뇌혈관수술비"):,}')
    if g('심장수술비'): r_set('TextBox 17',7,1,f': {g("심장수술비"):,}')
    if g('상해수술비'): r_set('TextBox 19',0,1,f': {g("상해수술비"):,}')
    if any(surg_s):
        r_set('TextBox 19',3,0,f'({"/".join(str(x) for x in surg_s)})')
        r_set('TextBox 19',3,2,'')
    if g('골절수술비'): r_set('TextBox 19',4,1,f': {g("골절수술비"):,}')

    # 실손
    실손_dates=[ct['contract_date'] for ct in contracts
        if any('실손' in k or '입원의료비' in k for k in ct['dambo']) and ct['contract_date']]
    실손가입일=min(실손_dates) if 실손_dates else '___________'
    by['TextBox 59'].text_frame.word_wrap=False
    by['TextBox 59'].text_frame.paragraphs[0].runs[0].text='실손'
    by['TextBox 59'].text_frame.paragraphs[1].runs[0].text='('
    by['TextBox 59'].text_frame.paragraphs[1].runs[1].text='가입일:'
    by['TextBox 59'].text_frame.paragraphs[1].runs[2].text=f'{실손가입일})'
    for r in by['TextBox 59'].text_frame.paragraphs[1].runs: r.font.size=Pt(8)
    if g('실손입원'): r_set('TextBox 6',0,1,f': {g("실손입원"):,}')
    if g('실손통원'): r_set('TextBox 6',1,1,f': {g("실손통원"):,} / ')
    if g('MRI'): r_set('TextBox 6',2,0,f'MRI : {g("MRI"):,}')
    if g('도수치료'): r_set('TextBox 6',3,1,f': {g("도수치료"):,}')
    if g('비급여주사'): r_set('TextBox 6',4,1,f': {g("비급여주사"):,}')

    # 상해/기타
    if g('골절진단비(치아제외)'): r_set('TextBox 7',0,1,f': {g("골절진단비(치아제외)"):,}')
    if g('화상진단비'): r_set('TextBox 7',2,1,f': {g("화상진단비"):,}')
    if g('깁스치료비'): r_set('TextBox 7',5,1,f': {g("깁스치료비"):,}')
    if g('응급실내원비(응급)'): r_set('TextBox 7',6,1,f': {g("응급실내원비(응급)"):,}')
    if g('일상생활배상책임'): r_set('TextBox 5',0,1,f': {g("일상생활배상책임"):,}')

    # 운전자
    if g('교통사고처리지원금'): r_set('TextBox 9',0,1,f': {g("교통사고처리지원금"):,}')
    if g('교통사고벌금(대물)'): r_set('TextBox 9',1,1,f': {g("교통사고벌금(대물)"):,}')
    if g('변호사선임비용'): r_set('TextBox 9',4,1,f': {g("변호사선임비용"):,}')

    # 입원/간병
    if g('질병일당'): r_set('TextBox 22',0,1,f': {g("질병일당"):,} / ')
    if g('상해일당'): r_set('TextBox 22',1,1,f': {g("상해일당"):,} / ')
    if g('1인실 상급병원일당'): r_set('TextBox 22',3,2,f': {g("1인실 상급병원일당"):,}')
    if g('1인실 종합병원일당'): r_set('TextBox 22',4,2,f': {g("1인실 종합병원일당"):,}')
    if g('간병인사용일당'): r_set('TextBox 22',7,1,f': {g("간병인사용일당"):,} / ')
    if g('요양병원일당'): r_set('TextBox 22',7,3,f': {g("요양병원일당"):,}')
    if g('간호통합병동일당'): r_set('TextBox 22',8,2,f': {g("간호통합병동일당"):,}')

    # 치아
    if g('크라운'): r_set('TextBox 13',0,1,f': {g("크라운"):,}')
    if g('임플란트'): r_set('TextBox 13',1,1,f': {g("임플란트"):,}')

    prs.save(path)
    return True

# ── 분석 요약 생성 ────────────────────────────────────────────────────
def make_summary(data):
    contracts=data['contracts']
    cust=data['client']
    total_premium=sum(ct['premium'] for ct in contracts)
    갱신수=sum(1 for ct in contracts if ct['renewal']=='갱신')
    비갱신수=len(contracts)-갱신수

    lines=[
        f"<b>👤 {cust} 고객님 분석 완료</b>",
        f"",
        f"📋 <b>계약 현황</b>",
        f"  • 총 계약 수: <b>{len(contracts)}건</b>",
        f"  • 갱신형: {갱신수}건 / 비갱신형: {비갱신수}건",
        f"  • 월 보험료 합계: <b>{total_premium:,}원</b>",
        f"",
        f"🏢 <b>가입 회사</b>",
    ]
    for ct in contracts:
        tag='🔵갱신' if ct['renewal']=='갱신' else '🔴비갱신' if '비갱신' in ct['renewal'] else '🟢완납'
        lines.append(f"  • {ct['company']} [{tag}] {ct['premium']:,}원")

    # 주요 담보 합계
    totals={}
    for ct in contracts:
        for raw,amt in ct['dambo'].items():
            std=resolve(raw)
            if std: totals[std]=totals.get(std,0)+amt

    key_items=[
        ('암진단비','🎗암진단'),('뇌혈관진단비','🧠뇌혈관'),('허혈성','❤️허혈성'),
        ('급성심근경색','❤️급성심근'),('상해사망','💀상해사망'),('질병사망(80세)','💀질병사망'),
        ('실손입원','🏥실손'),
    ]
    found=[(lbl,totals[k]) for k,lbl in key_items if k in totals and totals[k]>0]
    if found:
        lines.append(f"")
        lines.append(f"🔑 <b>주요 담보 합계 (만원)</b>")
        for lbl,amt in found:
            lines.append(f"  • {lbl}: <b>{amt:,}만원</b>")

    return '<br>'.join(lines)

# ── HTML ───────────────────────────────────────────────────────────────
INDEX_HTML = r'''<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0">
<title>MAKEONE 보장분석실</title>
<style>
  :root{--bg:#0c0d10;--panel:#15171c;--line:#2a2d34;--acc:#E0463B;--acc2:#F4897F;
    --ink:#EAECEF;--mute:#929aa6;--green:#4ADE80;--amber:#F5B547;--red:#FF6B6B;--blue:#5B9BFF}
  *{box-sizing:border-box;margin:0;padding:0;-webkit-tap-highlight-color:transparent}
  body{background:var(--bg);color:var(--ink);font-family:'Pretendard','Noto Sans KR',sans-serif;line-height:1.55}
  #gate{position:fixed;inset:0;z-index:100;background:var(--bg);display:flex;flex-direction:column;
    align-items:center;justify-content:center;padding:30px 26px;text-align:center}
  #gate .kick{font-size:14px;font-weight:800;letter-spacing:.45em;color:var(--acc);margin-bottom:14px}
  #gate h1{font-size:30px;font-weight:800;letter-spacing:-.01em;margin-bottom:14px}
  #gate .s{font-size:14px;color:var(--mute);margin-bottom:38px}
  #gate .pw{width:100%;max-width:420px;background:#1a1c22;border:1px solid var(--line);border-radius:14px;
    padding:18px 20px;font-size:17px;color:var(--ink);text-align:center;letter-spacing:.3em;outline:none}
  #gate .pw:focus{border-color:var(--acc)}
  #gate .go{width:100%;max-width:420px;margin-top:14px;border:none;border-radius:14px;padding:18px;
    font-size:17px;font-weight:800;color:#fff;background:var(--acc);cursor:pointer}
  #gate .go:active{transform:translateY(1px)}
  #gate .err{color:var(--acc2);font-size:13px;font-weight:700;margin-top:14px;min-height:18px}
  .shake{animation:sh .35s}
  @keyframes sh{0%,100%{transform:translateX(0)}25%{transform:translateX(-8px)}75%{transform:translateX(8px)}}
  .app{max-width:520px;margin:0 auto;height:100vh;display:none;flex-direction:column}
  header{padding:14px 18px;border-bottom:1px solid var(--line);
    background:linear-gradient(135deg,#1a1115,#0d0e11 60%,#1c1216);display:flex;align-items:center;gap:10px}
  .logo{width:32px;height:32px;border-radius:9px;border:1px solid var(--acc);display:flex;
    align-items:center;justify-content:center;font-size:16px;background:linear-gradient(135deg,#241012,#0d0e11)}
  h1{font-size:14px;font-weight:800}h1 b{color:var(--acc2)}
  .sub{font-size:10px;color:var(--mute)}
  .chat{flex:1;overflow-y:auto;padding:16px 12px;display:flex;flex-direction:column;gap:12px}
  .msg{max-width:90%;font-size:13px}
  .me{align-self:flex-end;background:rgba(224,70,59,.14);border:1px solid rgba(224,70,59,.32);
    border-radius:14px 14px 4px 14px;padding:9px 13px}
  .bot{align-self:flex-start;background:var(--panel);border:1px solid var(--line);
    border-radius:14px 14px 14px 4px;padding:11px 14px;width:100%}
  /* ★ 개별 파일 카드 */
  .file-cards{display:flex;flex-direction:column;gap:8px;margin-top:10px}
  .file-card{display:flex;align-items:center;gap:11px;border-radius:12px;
    padding:11px 13px;text-decoration:none;color:var(--ink)}
  .file-card.xl{background:rgba(74,222,128,.06);border:1px solid rgba(74,222,128,.3)}
  .file-card.pt{background:rgba(91,155,255,.06);border:1px solid rgba(91,155,255,.3)}
  .file-card .ic{font-size:22px}.file-card .nm{flex:1;font-size:12.5px;font-weight:700}
  .file-card .dl{font-size:11px;font-weight:800;padding:5px 11px;border-radius:8px}
  .file-card.xl .dl{color:var(--green);background:rgba(74,222,128,.12)}
  .file-card.pt .dl{color:var(--blue);background:rgba(91,155,255,.12)}
  .summary-box{background:#1a1f2a;border:1px solid #2a3040;border-radius:10px;
    padding:12px 14px;margin-top:10px;font-size:12px;line-height:1.7;white-space:pre-wrap}
  .err{color:#ffb4b4;font-size:12px}
  .spin{width:22px;height:22px;border:3px solid var(--line);border-top-color:var(--acc);
    border-radius:50%;animation:sp .8s linear infinite;display:inline-block;vertical-align:middle}
  @keyframes sp{to{transform:rotate(360deg)}}
  .bar{padding:12px;border-top:1px solid var(--line);display:flex;gap:9px;background:var(--bg)}
  .up{flex:1;border:1.5px dashed rgba(224,70,59,.5);border-radius:12px;padding:13px;text-align:center;
    font-size:13px;font-weight:700;cursor:pointer;color:var(--acc2)}
  .send{border:none;border-radius:12px;padding:0 20px;font-weight:800;font-size:14px;background:var(--acc);color:#fff;cursor:pointer}
  .send:disabled{opacity:.4}
  footer{text-align:center;font-size:10px;color:var(--mute);padding:8px}
  footer b{color:var(--acc2)}
</style></head>
<body>
<div id="gate">
  <div class="kick">MAKEONE</div>
  <h1>MAKEONE 보장분석실</h1>
  <div class="s">접속 비밀번호를 입력하세요</div>
  <input id="pw" class="pw" type="password" inputmode="numeric" placeholder="비밀번호" autocomplete="off">
  <button id="go" class="go">접속</button>
  <div id="gerr" class="err"></div>
</div>
<div class="app" id="app">
  <header><div class="logo">📋</div><div><h1>MAKEONE <b>보장분석실</b></h1>
    <div class="sub">TXT 넣으면 → 엑셀+PPT 개별 다운로드 · 최은혜 지점장</div></div></header>
  <div class="chat" id="chat">
    <div class="msg bot">
      보장분석지 <b>TXT 파일</b>을 올려주세요. 엑셀·PPT 파일을 각각 드려요.<br><br>
      <span style="font-size:11px;color:var(--mute)">
        💡 Adobe Acrobat → 편집 → 모두선택(Ctrl+A) → 복사(Ctrl+C)<br>
        → 메모장 붙여넣기 → .txt 저장 → 여기 업로드
      </span>
    </div>
  </div>
  <div class="bar">
    <label class="up" id="up">📄 <span id="uplabel">보장분석지 TXT 선택</span></label>
    <button class="send" id="send" disabled>분석</button>
  </div>
  <footer>미래를 <b>바르게</b> 설계합니다 · BARUM</footer>
</div>
<input type="file" id="fi" accept=".txt,text/plain" style="display:none">
<script>
const $=s=>document.querySelector(s);
let ACCESS='';
async function unlock(){
  const v=$("#pw").value;
  $("#gerr").textContent="확인 중…";
  try{
    const r=await fetch("/check",{method:"POST",headers:{"Content-Type":"application/json"},body:JSON.stringify({pw:v})});
    const j=await r.json();
    if(j.ok){ACCESS=v;$("#gerr").textContent="";$("#gate").style.display="none";$("#app").style.display="flex";}
    else{fail();}
  }catch(e){$("#gerr").textContent="서버 연결 실패";}
}
function fail(){
  $("#gerr").textContent="비밀번호가 올바르지 않습니다.";
  $("#gate").classList.add("shake");
  setTimeout(()=>$("#gate").classList.remove("shake"),350);
  $("#pw").value="";$("#pw").focus();
}
$("#go").onclick=unlock;
$("#pw").addEventListener("keydown",e=>{if(e.key==="Enter")unlock();});
window.addEventListener("load",()=>$("#pw").focus());

const chat=$("#chat");let file=null;
$("#up").onclick=()=>$("#fi").click();
$("#fi").onchange=e=>{
  file=e.target.files[0];
  if(file){$("#uplabel").textContent=file.name;$("#send").disabled=false;}
};
function esc(s){return String(s==null?"":s).replace(/[&<>]/g,c=>({"&":"&amp;","<":"&lt;",">":"&gt;"}[c]));}
function add(html,cls){const d=document.createElement("div");d.className="msg "+cls;d.innerHTML=html;chat.appendChild(d);chat.scrollTop=chat.scrollHeight;return d;}

function b64toBlob(b64,mime){
  const bin=atob(b64);const arr=new Uint8Array(bin.length);
  for(let i=0;i<bin.length;i++)arr[i]=bin.charCodeAt(i);
  return new Blob([arr],{type:mime});
}
function triggerDownload(blob,fname){
  const u=URL.createObjectURL(blob);
  const a=document.createElement("a");a.href=u;a.download=fname;a.click();
  setTimeout(()=>URL.revokeObjectURL(u),3000);
}

$("#send").onclick=async()=>{
  if(!file)return;
  add("📄 "+esc(file.name),"me");
  $("#send").disabled=true;$("#up").style.opacity=.5;
  const loading=add(
    '<div style="display:flex;align-items:center;gap:11px"><span class="spin"></span>'+
    '<div style="flex:1"><div id="ldmsg" style="font-weight:800">📄 TXT 파싱 중…</div>'+
    '<div id="ldtime" style="font-size:11px;color:var(--mute);margin-top:2px">0초 · 기다려 주세요</div></div></div>',"bot");
  const t0=Date.now();
  const steps=["📄 TXT 파싱 중…","🔎 별첨 담보 추출 중…","📊 엑셀 생성 중…","🖼 PPT 채우는 중…","✅ 완성 중…"];
  let si=0;
  const timer=setInterval(()=>{
    si=Math.min(si+1,steps.length-1);
    const s=Math.floor((Date.now()-t0)/1000);
    const tm=document.getElementById("ldtime");
    const mm=document.getElementById("ldmsg");
    if(tm)tm.textContent=s+"초 경과";
    if(mm)mm.textContent=steps[si];
  },8000);
  const fd=new FormData();fd.append("file",file);fd.append("pw",ACCESS);
  try{
    const r=await fetch("/analyze",{method:"POST",body:fd});
    clearInterval(timer);loading.remove();
    const j=await r.json();
    if(!j.ok){add('<span class="err">⚠ '+esc(j.error||"실패")+'</span>',"bot");}
    else{
      // ★ 엑셀 + PPT 개별 다운로드
      const xlBlob=b64toBlob(j.xlsx_b64,"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet");
      triggerDownload(xlBlob,j.xlsx_name);
      let ptCard='';
      if(j.pptx_b64){
        const ptBlob=b64toBlob(j.pptx_b64,"application/vnd.openxmlformats-officedocument.presentationml.presentation");
        setTimeout(()=>triggerDownload(ptBlob,j.pptx_name),800);
        ptCard=`<a class="file-card pt" href="#" onclick="event.preventDefault()">
          <span class="ic">📊</span>
          <span class="nm">${esc(j.pptx_name)}<br><span style="font-size:10px;color:var(--mute)">보장분석 PPT</span></span>
          <span class="dl">자동저장 ⬇</span></a>`;
      }
      // ★ 분석 요약 + 파일 카드
      add(
        '<b>✅ 분석 완료!</b>'+
        '<div class="summary-box">'+j.summary+'</div>'+
        '<div class="file-cards">'+
        `<a class="file-card xl" href="#" onclick="event.preventDefault()">
          <span class="ic">📗</span>
          <span class="nm">${esc(j.xlsx_name)}<br><span style="font-size:10px;color:var(--mute)">보장진단 엑셀</span></span>
          <span class="dl">자동저장 ⬇</span></a>`+
        ptCard+
        '</div>',"bot");
    }
  }catch(e){clearInterval(timer);loading.remove();add('<span class="err">오류: '+esc(e.message)+'</span>',"bot");}
  file=null;$("#uplabel").textContent="보장분석지 TXT 선택";
  $("#send").disabled=true;$("#fi").value="";$("#up").style.opacity=1;
};
</script></body></html>
'''

# ── FastAPI ────────────────────────────────────────────────────────────
app=FastAPI()
ACCESS_PW=os.environ.get('ACCESS_PW','1009')

@app.get('/health')
def health(): return {'ok':True,'version':'v6-split-download'}

@app.get('/',response_class=HTMLResponse)
def home(): return INDEX_HTML

@app.post('/check')
async def check_pw(body:dict): return {'ok':body.get('pw')==ACCESS_PW}

@app.post('/analyze')
async def analyze(file:UploadFile=File(...),pw:str=Form('')):
    if pw!=ACCESS_PW: return JSONResponse({'ok':False,'error':'비밀번호 오류'})
    if not (file.filename.lower().endswith('.txt') or 'text' in (file.content_type or '')):
        return JSONResponse({'ok':False,'error':'TXT 파일만 가능. Adobe Acrobat 텍스트 추출본을 업로드하세요.'})
    raw=await file.read()
    for enc in ['utf-8','cp949','euc-kr']:
        try: txt=raw.decode(enc); break
        except: pass
    else: txt=raw.decode('utf-8',errors='ignore')

    try:
        data=parse_txt(txt)
        if not data['contracts']:
            return JSONResponse({'ok':False,'error':'계약을 찾지 못했습니다. (정상계약 리스트) 섹션이 포함된 TXT인지 확인하세요.'})
        cust=data['client']
        d=tempfile.mkdtemp()
        now=datetime.datetime.now()
        xl=os.path.join(d,f'보장진단_{cust}.xlsx')
        pt=os.path.join(d,f'보장분석지_{cust}.pptx')
        build_excel(data,xl)
        ppt_ok=build_ppt(data,pt)

        # ★ 엑셀 + PPT 개별 base64 반환 (ZIP 없음)
        xlsx_b64=base64.b64encode(open(xl,'rb').read()).decode()
        xlsx_name=f'보장진단_{cust}.xlsx'
        response={'ok':True,'xlsx_b64':xlsx_b64,'xlsx_name':xlsx_name,
                  'summary':make_summary(data),'pptx_ready':ppt_ok}
        if ppt_ok and os.path.exists(pt):
            response['pptx_b64']=base64.b64encode(open(pt,'rb').read()).decode()
            response['pptx_name']=f'보장분석지_{cust}.pptx'
        return JSONResponse(response)
    except Exception as e:
        return JSONResponse({'ok':False,'error':str(e),'trace':traceback.format_exc()[-1500:]})
