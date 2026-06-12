# -*- coding: utf-8 -*-
import re
from pptx import Presentation

SRC='/mnt/user-data/uploads/2025_10_보장분석지.pptx'
OUT='/home/claude/work/pkg/MASTER_blank.pptx'

AMOUNT = re.compile(r'\d\s*[천만억원]|\d+\s*[+/]\s*\d')  # 값 표식: 2천만 / 500+500 / 20/20
CLEAR_BOX = {'TextBox 1','TextBox 16','TextBox 24','TextBox 25'}  # 이름·날짜 = 항상 비움

def keepmask(S):
    """문단 전체 텍스트 기준, 지울 문자 인덱스 판정. 라벨·콜론·슬래시(필드구분)는 keep."""
    if ':' not in S:
        # 콜론 없음: 값표식 있으면 값박스로 보고 전부 삭제, 아니면 라벨이라 보존
        return [not bool(AMOUNT.search(S))]*len(S)
    keep=[True]*len(S); mode='L'; i=0
    while i<len(S):
        c=S[i]
        if mode=='L':
            keep[i]=True
            if c==':': mode='V'
            i+=1
        else:  # 값 모드: 값 문자 삭제
            if c=='/':
                rest=S[i+1:]
                if ':' in rest:        # 뒤에 또 라벨 → 필드구분자, '/' 유지하고 라벨 모드
                    keep[i]=True; mode='L'
                else:                   # 값 대안(2천/1천) → '/' 삭제
                    keep[i]=False
                i+=1
            elif c==' ':
                keep[i]=True; i+=1      # 공백은 유지(라벨 : 모양 보존)
            else:
                keep[i]=False; i+=1
    return keep

JONG = re.compile(r'^\(\s*\d+(\s*/\s*\d+)+\s*\)$')  # ( 20/20/100/200/500 ) 종별슬롯

def blank_para(p):
    S=''.join(r.text for r in p.runs)
    if not S.strip(): return None
    # 1~5종 종별 슬롯: 숫자만 비우고 ( / / / / ) 골격 유지
    if JONG.match(S.strip()):
        n=len(re.findall(r'\d+', S))
        skel='(  '+' / '.join(['   ']*n)+'  )'
        runs=p.runs
        if runs:
            runs[0].text=skel
            for r in runs[1:]: r.text=''
        return [(S, skel)]
    km=keepmask(S)
    out=[]; idx=0
    for r in p.runs:
        n=len(r.text)
        new=''.join(ch for j,ch in enumerate(r.text) if km[idx+j])
        if new!=r.text: out.append((r.text,new))
        r.text=new
        idx+=n
    return out

prs=Presentation(SRC)
# slide8만 남긴다 (rel까지 정리)
from pptx.oxml.ns import qn
lst=prs.slides._sldIdLst
for i,sid in enumerate(list(lst)):
    if i!=8:
        rId=sid.get(qn('r:id'))
        try: prs.part.drop_rel(rId)
        except Exception: pass
        lst.remove(sid)
s=prs.slides[0]

changes=[]
for sh in s.shapes:
    if not sh.has_text_frame: continue
    tf=sh.text_frame
    if sh.name in CLEAR_BOX:
        for p in tf.paragraphs:
            for r in p.runs:
                if r.text.strip(): changes.append((sh.name,r.text,'')); r.text=''
        continue
    for p in tf.paragraphs:
        ch=blank_para(p)
        if ch:
            for a,b in ch: changes.append((sh.name,a,b))

prs.save(OUT)
print('SAVED',OUT,'| 변경 run',len(changes))
print('--- 삭제/변경된 값 (앞→뒤) ---')
for nm,a,b in changes:
    if a.strip()!=b.strip():
        print(f'  [{nm}] "{a}" -> "{b}"')
